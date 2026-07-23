"""Build the 5 remaining new slides into a temp deck for splicing."""
import sys
from pathlib import Path
sys.path.insert(0, str(Path(__file__).parent))
import _build_Italy_Class1 as B  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE  # noqa: E402

HERE = Path(__file__).parent
IMG = HERE / "Images"
GREY = "E7E9EC"


def takeaway(slide, text):
    y, h = int(Inches(6.5)), int(Inches(0.52))
    r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, int(B.MARGIN), y, int(B.RULE_W), h)
    r.fill.solid(); r.fill.fore_color.rgb = B.GOLD; r.line.fill.background()
    r.shadow.inherit = False
    try:
        r.adjustments[0] = 0.25
    except Exception:
        pass
    _shadow(r)
    B._add_text(slide, int(B.MARGIN), y, int(B.RULE_W), h, text, size=18, bold=True,
                italic=True, color=B.NAVY, font="Calibri", align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)


def _shadow(shape):
    from lxml import etree as ET
    from pptx.oxml.ns import qn
    spPr = shape._element.spPr
    for e in spPr.findall(qn('a:effectLst')):
        spPr.remove(e)
    lst = ET.SubElement(spPr, qn('a:effectLst'))
    sh = ET.SubElement(lst, qn('a:outerShdw'))
    sh.set('blurRad', '50800'); sh.set('dist', '38100'); sh.set('dir', '2700000'); sh.set('rotWithShape', '0')
    c = ET.SubElement(sh, qn('a:srgbClr')); c.set('val', '000000')
    ET.SubElement(c, qn('a:alpha')).set('val', '30000')


def img_caption(slide, text, x, y, w):
    B._add_text(slide, int(x), int(y), int(w), int(Inches(0.3)), text, size=11,
                italic=True, color=B.GRAY, font="Calibri", align=PP_ALIGN.CENTER)


def content_img(prs, section, title, items, img, cap, cx, cy, mw, mh, capx, capy, capw):
    s = B._blank_slide(prs)
    B._top_bar(s, section)
    B._action_title(s, title)
    B._hbullets(s, items, left=int(B.MARGIN), top=int(Inches(1.9)),
                width=int(Inches(7.1)), height=int(Inches(4.2)))
    B._place_image(s, img, cx=int(cx), cy=int(cy), max_w=int(mw), max_h=int(mh),
                   shadow=True, rounded=True)
    img_caption(s, cap, capx, capy, capw)
    return s


def s_renaissance(prs):
    s = content_img(prs, "The Renaissance",
        "The Renaissance Economy: Italy Invents Modern Business", [
            (0, "Double-entry bookkeeping (codified by Pacioli, 1494) — still used worldwide"),
            (0, "The gold florin (from 1252): an international reserve currency"),
            (0, "Merchant banks — Medici, Bardi, Peruzzi — financed kings and popes"),
            (0, "Bills of exchange, marine insurance, venture partnerships (colleganza)"),
            (0, "Venice’s Arsenal: a proto-assembly-line shipyard"),
        ], IMG / "florin_crop.jpg",
        "Gold florin of Florence (PAS, CC BY 2.0)",
        cx=Inches(10.55), cy=Inches(3.35), mw=Inches(4.4), mh=Inches(2.4),
        capx=Inches(8.3), capy=Inches(4.7), capw=Inches(4.5))
    takeaway(s, "Italy invented the toolkit of modern capitalism — banking, accounting, insurance, the firm")
    B._footer(s, 43)


def s_industry(prs):
    s = content_img(prs, "Industry & the State",
        "Late Industrialization and the Rise of the State", [
            (0, "Industry came late (from the 1890s), concentrated in the North: Milan–Turin–Genoa"),
            (0, "Fiat (1899), Pirelli (1872), Olivetti (1908) — powered by Alpine hydro"),
            (0, "Financed by German-style universal banks (Banca Commerciale, 1894)"),
            (0, "1933: the state creates IRI — a giant holding company (steel, ships, telecoms)"),
            (0, "IRI’s public-enterprise model shaped Italian capitalism into the 1990s"),
        ], IMG / "lingotto.jpg",
        "Fiat Lingotto factory, Turin (J.-P. Dalbéra, CC BY 2.0)",
        cx=Inches(10.5), cy=Inches(3.35), mw=Inches(4.6), mh=Inches(2.9),
        capx=Inches(8.1), capy=Inches(5.0), capw=Inches(4.9))
    takeaway(s, "Italy industrialized late, in the North, via banks — and then the state")
    B._footer(s, 63)


def s_miracle(prs):
    s = content_img(prs, "Post-WWII & Marshall Plan",
        "The Economic Miracle (1958–63)", [
            (0, "Growth on the order of 6–8% a year — agrarian to industrial in a generation"),
            (0, "Mass migration from the South to the northern factories of Turin and Milan"),
            (0, "Icons: the Vespa (1946), the Fiat 500 (1957), household “white goods”"),
            (0, "ENI (Enrico Mattei) builds a national energy champion; Autostrada del Sole (1964)"),
        ], IMG / "fiat500.jpg",
        "Fiat Nuova 500, 1957 (T. Doerfer, CC BY 3.0)",
        cx=Inches(10.5), cy=Inches(3.4), mw=Inches(4.7), mh=Inches(3.1),
        capx=Inches(8.1), capy=Inches(5.15), capw=Inches(4.9))
    takeaway(s, "Export manufacturing plus management catch-up transformed Italy — fast")
    B._footer(s, 89)


def s_thirditaly(prs):
    s = content_img(prs, "The Italian Economy Today",
        "The “Third Italy”: Industrial Districts & Made in Italy", [
            (0, "From the 1970s: clusters of small, family-owned firms in the Center & North-East"),
            (0, "Prato textiles, Sassuolo ceramics, Belluno eyewear (Luxottica), Bologna packaging"),
            (0, "Flexible specialization + local know-how + export niches = “Made in Italy”"),
            (0, "Rooted in the medieval commune/guild tradition — and in family capitalism"),
        ], IMG / "murano.jpg",
        "Murano glassmaker, Venice (Saffron Blaze, CC BY-SA 3.0)",
        cx=Inches(10.5), cy=Inches(3.45), mw=Inches(4.7), mh=Inches(3.2),
        capx=Inches(8.1), capy=Inches(5.25), capw=Inches(4.9))
    takeaway(s, "Italy competes on design, quality, and niche manufacturing — by networked family SMEs")
    B._footer(s, 90)


def _card(slide, x, y, w, h, fill, header, items):
    from pptx.dml.color import RGBColor
    r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, int(x), int(y), int(w), int(h))
    r.fill.solid(); r.fill.fore_color.rgb = RGBColor.from_string(fill)
    r.line.color.rgb = B.GOLD; r.line.width = Pt(1.25)
    r.shadow.inherit = False
    try:
        r.adjustments[0] = 0.06
    except Exception:
        pass
    _shadow(r)
    B._add_text(slide, int(x), int(y + Inches(0.12)), int(w), int(Inches(0.5)), header,
                size=20, bold=True, color=B.NAVY, font="Calibri", align=PP_ALIGN.CENTER)
    B._hbullets(slide, [(0, t) for t in items], left=int(x + Inches(0.3)),
                top=int(y + Inches(0.7)), width=int(w - Inches(0.6)),
                height=int(h - Inches(0.85)), anchor=MSO_ANCHOR.TOP,
                main_space=8)


def s_paradox(prs):
    s = B._blank_slide(prs)
    B._top_bar(s, "The Italian Economy Today")
    B._action_title(s, "The Paradox: World-Class Brands, Stalled Productivity")
    _card(s, Inches(0.35), Inches(1.75), Inches(6.2), Inches(4.4), "FDF6E6", "Strengths", [
        "Design & craftsmanship",
        "Global brands: fashion, food, luxury cars, machinery",
        "Family firms & “pocket multinationals”",
        "Industrial districts (“Made in Italy”)",
    ])
    _card(s, Inches(6.75), Inches(1.75), Inches(6.2), Inches(4.4), GREY, "Constraints", [
        "Flat productivity since the mid-1990s",
        "Small firm size; weak meritocracy",
        "The North–South divide (~2:1)",
        "Demographics & brain drain; weak institutions",
    ])
    takeaway(s, "Excellence is sectoral and relational; scale and institutions are the persistent weak spots")
    B._footer(s, 98)


def main():
    prs = Presentation()
    prs.slide_width = B.SLIDE_W
    prs.slide_height = B.SLIDE_H
    s_renaissance(prs)
    s_industry(prs)
    s_miracle(prs)
    s_thirditaly(prs)
    s_paradox(prs)
    prs.save(HERE / "_moreslides.pptx")
    print("wrote _moreslides.pptx (5 slides)")


if __name__ == "__main__":
    main()
