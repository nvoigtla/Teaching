# -*- coding: utf-8 -*-
"""Build "Module 1 - Example Candidates.pptx" — a review deck of recent
(2023–2026) real-world examples for the Module 1 concepts, from which
Nico picks slides to adopt into "Module 1 - Revised.pptx".

One slide per candidate: concept tag in the top bar, takeaway title,
fact bullets (all facts verified against the listed sources), a
teaching-angle line, a discussion question, and a source line. Candidates
are text-first by design — the adopted ones get the full visual
treatment (images/charts) when ported into the main deck.

Reuses the Module 1 helper layer (import side effects: defs only).
"""
from _build_Module1 import (
    CREAM, FADED, GOLD, GOLD_W, GRAY, MARGIN, NAVY, RED, RULE, RULE_W,
    SLIDE_H, SLIDE_W, WHITE, STEEL, GREEN_BR, GREEN_MB, BLUE_PED,
    Inches, Pt, PP_ALIGN, MSO_ANCHOR, MSO_SHAPE, Presentation, Path,
    _add_hierarchical_bullets, _add_rect, _add_text, _add_slidenum_field,
    _blank_slide, _draw_action_title, _draw_top_bar_tc, _add_drop_shadow,
    _add_rounded_filled_box, _add_outlined_box, _add_convention_box,
    _add_media_image, _add_arrow, _add_arrow_shape, _set_notes,
    SimpleFig, _fig_axes, _fig_line, _fig_guide, _fig_xlab, _fig_ylab,
    _fig_curve_label,
)

OUT_DIR = Path(__file__).parent
FOOTER_TEXT = "Management 405  ·  Module 1  ·  Example Candidates (for review)"


def _footer(slide, page_num):
    _add_rect(slide, 0, Inches(7.15), SLIDE_W, Inches(0.02), RULE)
    _add_rect(slide, MARGIN, Inches(7.135), GOLD_W, Inches(0.05), GOLD)
    _add_text(slide, MARGIN, Inches(7.20), Inches(11), Inches(0.32),
              FOOTER_TEXT, size=12, color=GRAY)
    _add_slidenum_field(slide, Inches(12.55), Inches(7.20), Inches(0.55),
                        Inches(0.32), page_num)


def slide_cover(prs):
    slide = _blank_slide(prs)
    _add_text(slide, Inches(0.9), Inches(2.0), SLIDE_W - Inches(1.8),
              Inches(1.1), "Example Candidates",
              size=54, bold=True, color=NAVY, font="Calibri",
              align=PP_ALIGN.CENTER)
    _add_text(slide, 0, Inches(3.35), SLIDE_W, Inches(0.7),
              "Module 1 · recent real-world examples for review",
              size=28, bold=True, color=GOLD, font="Calibri",
              align=PP_ALIGN.CENTER)
    _add_rect(slide, int((SLIDE_W - Inches(4.0)) / 2), Inches(4.35),
              Inches(4.0), 54864, GOLD)
    _add_text(slide, Inches(1.7), Inches(4.75), SLIDE_W - Inches(3.4),
              Inches(1.5),
              "One slide per candidate. Facts verified against the sources "
              "listed on each slide (retrieved 2026-08-20). Adopted "
              "candidates get the full visual treatment (images, native "
              "charts) when ported into the main deck.",
              size=16, italic=True, color=GRAY, font="Calibri",
              align=PP_ALIGN.CENTER)
    _add_rect(slide, 0, Inches(7.15), SLIDE_W, Inches(0.02), RULE)
    _add_rect(slide, MARGIN, Inches(7.135), GOLD_W, Inches(0.05), GOLD)
    return slide


def candidate_slide(prs, page_num, concept, title, facts, angle,
                    discussion, sources, *, visual=None, notes=None):
    """One candidate: concept in the tag, takeaway title, fact bullets,
    cream teaching-angle box, discussion line, gray source line."""
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, "Module 1 · Candidates · " + concept)
    if len(title) > 62:
        _add_text(slide, MARGIN, Inches(0.55), RULE_W, Inches(0.7),
                  title, size=26, bold=True, color=NAVY, font="Calibri")
        _add_rect(slide, MARGIN, Inches(1.25), RULE_W, Inches(0.02), RULE)
        _add_rect(slide, MARGIN, Inches(1.235), GOLD_W, Inches(0.05), GOLD)
    else:
        _draw_action_title(slide, title)

    items = [(f, 0) if isinstance(f, str) else f for f in facts]
    box = _add_hierarchical_bullets(
        slide, left=MARGIN, top=Inches(1.55), width=RULE_W,
        height=Inches(3.0), items=items,
        size=20, sub_size=18, line_spacing_pts=8)
    box.text_frame.vertical_anchor = MSO_ANCHOR.TOP

    # teaching-angle box (cream card) — all card text >= 18 pt
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, int(MARGIN), int(Inches(4.55)),
        int(RULE_W), int(Inches(1.38)))
    try:
        card.adjustments[0] = 0.10
    except Exception:
        pass
    card.fill.solid(); card.fill.fore_color.rgb = CREAM
    card.line.color.rgb = NAVY; card.line.width = Pt(1.0)
    card.shadow.inherit = False
    _add_drop_shadow(card)
    tb = slide.shapes.add_textbox(int(MARGIN + Inches(0.2)),
                                  int(Inches(4.62)),
                                  int(RULE_W - Inches(0.4)),
                                  int(Inches(1.24)))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = "Teaching angle:  "
    r1.font.name = "Calibri"; r1.font.size = Pt(18)
    r1.font.bold = True; r1.font.color.rgb = NAVY
    r2 = p.add_run(); r2.text = angle
    r2.font.name = "Calibri"; r2.font.size = Pt(18)
    r2.font.color.rgb = NAVY
    if visual:
        p2 = tf.add_paragraph()
        v1 = p2.add_run(); v1.text = "Proposed visual:  "
        v1.font.name = "Calibri"; v1.font.size = Pt(18)
        v1.font.bold = True; v1.font.color.rgb = GRAY
        v2 = p2.add_run(); v2.text = visual
        v2.font.name = "Calibri"; v2.font.size = Pt(18)
        v2.font.italic = True; v2.font.color.rgb = GRAY

    # discussion line
    db = slide.shapes.add_textbox(int(MARGIN), int(Inches(6.02)),
                                  int(RULE_W), int(Inches(0.62)))
    dtf = db.text_frame
    dtf.word_wrap = True
    dp = dtf.paragraphs[0]
    d1 = dp.add_run(); d1.text = "Discussion:  "
    d1.font.name = "Calibri"; d1.font.size = Pt(18)
    d1.font.bold = True; d1.font.color.rgb = GOLD
    d2 = dp.add_run(); d2.text = discussion
    d2.font.name = "Calibri"; d2.font.size = Pt(18)
    d2.font.italic = True; d2.font.color.rgb = NAVY

    # source line
    _add_text(slide, MARGIN, Inches(6.72), RULE_W, Inches(0.4),
              "Sources: " + sources, size=11, italic=True, color=GRAY,
              font="Calibri")

    _footer(slide, page_num)
    if notes:
        _set_notes(slide, notes)
    return slide


# --------------------------------------------------------------------------
# Shared bits for the EXPANDED (adoption-ready) example slides
# --------------------------------------------------------------------------

def _photo_caption(slide, left, top, width, text="Photos: Wikimedia Commons"):
    return _add_text(slide, left, top, width, Inches(0.28), text,
                     size=11, italic=True, color=GRAY, font="Calibri",
                     align=PP_ALIGN.CENTER)


def _quote_box(slide, left, top, width, height, quote, attribution, *,
               size=18):
    return _add_convention_box(
        slide, left, top, width, height,
        runs=[(quote, {'italic': True, 'size': size}),
              ("   — " + attribution,
               {'bold': True, 'size': max(size - 2, 16),
                'newline': True})],
        align=PP_ALIGN.LEFT)


def _source_line(slide, text, *, top=Inches(6.78)):
    return _add_text(slide, MARGIN, top, RULE_W, Inches(0.3),
                     "Sources: " + text, size=11, italic=True, color=GRAY,
                     font="Calibri")


# --------------------------------------------------------------------------
# EXPANDED 1a/1b — Tapestry–Capri (market definition)
# --------------------------------------------------------------------------

def exp_tapestry_case(prs, page_num):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, "Module 1 · Candidates · Market Definition")
    _draw_action_title(slide, "Market Definition Mini-Case: Tapestry–Capri")
    box = _add_hierarchical_bullets(
        slide, left=MARGIN, top=Inches(1.50), width=Inches(8.0),
        height=Inches(5.3),
        items=[
            ("Aug 2023: Tapestry (Coach, Kate Spade) agrees to buy Capri "
             "(Michael Kors, Versace) for $8.5B", 0),
            ([("The definition of the market", {'bold': True}),
              (" would turn out to be crucial for the case:", {})], 0, {}),
            ([("FTC: the market is ", {}),
              ("“accessible luxury” handbags", {'color': RED}),
              (" — roughly $100 to under $1,000", {})], 1, {}),
            ([("The firms: the market is ", {}),
              ("all handbags", {'color': RED}),
              (" — from fast fashion to Hermès — and entry is easy",
               {})], 1, {}),
        ],
        size=24, sub_size=22, line_spacing_pts=16,
        sub_line_spacing_pts=8)
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    _add_media_image(slide, "web_coach.jpg",
                     left=Inches(8.55), top=Inches(1.70),
                     width=Inches(3.67), rounded=True)     # h ≈ 2.45
    _add_media_image(slide, "web_michaelkors.jpg",
                     left=Inches(8.75), top=Inches(4.30),
                     width=Inches(3.27), rounded=True)     # h ≈ 2.45
    _photo_caption(slide, Inches(8.55), Inches(6.80), Inches(3.67))
    _footer(slide, page_num)
    _set_notes(slide, (
        "Expanded from the candidate list (research verified 2026-08-20). "
        "Preliminary injunction: Judge Jennifer Rochon, S.D.N.Y., 24 Oct "
        "2024 (169-page opinion); FTC administrative suit filed 22 Apr "
        "2024; merger agreement terminated 13 Nov 2024. The defense "
        "called the FTC's market 'gerrymandered'; the court found the "
        "'accessible luxury' segment real — distinct prices, customers, "
        "discounting, outlet distribution.\n"
        "Photos: Wikimedia Commons (Coach store, Tenmaya Fukuyama; "
        "Michael Kors store, Rehoboth Beach DE)."))
    return slide


def exp_tapestry_evidence(prs, page_num):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, "Module 1 · Candidates · Market Definition")
    _draw_action_title(slide,
                       "The Firms' Own Documents Drew the Market Boundary")
    # price-tier ladder
    y0, h0 = Inches(1.55), Inches(1.30)
    _add_rounded_filled_box(slide, Inches(0.75), y0, Inches(3.55), h0,
                            "Mass market\nunder $100",
                            fill=FADED, text_color=WHITE, size=19)
    _add_rounded_filled_box(
        slide, Inches(4.55), y0, Inches(4.25), h0,
        "“Accessible luxury”\n$100 – under $1,000\nCoach · Kate Spade · "
        "Michael Kors",
        fill=GOLD, text_color=NAVY, size=18)
    _add_rounded_filled_box(slide, Inches(9.05), y0, Inches(3.55), h0,
                            "True luxury\n$1,000+ · LV · Hermès · Chanel",
                            fill=FADED, text_color=WHITE, size=19)
    _add_text(slide, MARGIN, Inches(2.93), RULE_W, Inches(0.32),
              "the FTC's market — the firms' own term in SEC filings and "
              "investor decks, until the FTC sued",
              size=18, bold=True, italic=True,
              color=GOLD, font="Calibri", align=PP_ALIGN.CENTER)

    _add_text(slide, MARGIN, Inches(3.32), RULE_W, Inches(0.35),
              "Combined Tapestry + Capri share of “accessible luxury” "
              "handbags:", size=18, bold=True, color=NAVY, font="Calibri",
              align=PP_ALIGN.CENTER)
    cards = [("59%", "FTC's expert (third-party data)"),
             ("77%", "Capri's internal documents"),
             ("83%", "Tapestry's internal data")]
    for i, (num, lab) in enumerate(cards):
        x = Inches(1.35 + i * 3.65)
        shp = _add_outlined_box(slide, x, Inches(3.72), Inches(3.35),
                                Inches(1.10), "", rounded=True, shadow=True,
                                line=GOLD, line_w=1.5)
        tf = shp.text_frame
        p = tf.paragraphs[0]
        r1 = p.add_run(); r1.text = num + "   "
        r1.font.name = "Calibri"; r1.font.size = Pt(30)
        r1.font.bold = True; r1.font.color.rgb = NAVY
        r2 = p.add_run(); r2.text = lab
        r2.font.name = "Calibri"; r2.font.size = Pt(18)
        r2.font.color.rgb = GRAY

    _add_text(slide, MARGIN, Inches(4.92), RULE_W, Inches(0.30),
              "(figures from documents the companies had to hand over in "
              "the merger review — not leaked)", size=18, italic=True,
              color=GRAY, font="Calibri", align=PP_ALIGN.CENTER)
    _quote_box(slide, MARGIN + Inches(0.35), Inches(5.34),
               RULE_W - Inches(0.7), Inches(1.00),
               "“Bottom line, saying we're in the same market with true "
               "luxury is a joke. … Nobody says ‘should I buy a LV bag or "
               "a Coach bag?’”",
               "internal Tapestry message cited by the court")
    # main takeaway: gold bar, big navy bold (sources live in the notes)
    _add_rounded_filled_box(
        slide, Inches(0.90), Inches(6.44), Inches(11.53), Inches(0.60),
        "Oct 2024: the court sides with the FTC and blocks the deal — "
        "merger abandoned Nov 2024",
        fill=GOLD, text_color=NAVY, size=19, bold=True, corner_pct=0.18)
    _footer(slide, page_num)
    _set_notes(slide, (
        "Share figures (per the Clifford Chance briefing citing the "
        "Opinion at 97, verified 2026-08-20): 58.7% per the FTC's "
        "economic expert 'from largely third-party data'; 77% calculated "
        "in Capri's internal documents; 83% calculated from Tapestry's "
        "internal data. Provenance: ordinary-course internal documents "
        "that 'come to light during a merger investigation' — compulsory "
        "production, not leaks. The term 'accessible luxury' was the "
        "firms' own: used extensively in SEC filings and investor "
        "presentations before and after the acquisition, 'only for the "
        "term to disappear from their lexicon' after the FTC sued (they "
        "then preferred 'expressive luxury'). On-slide source line "
        "dropped for space — cite: Judge Rochon opinion, S.D.N.Y. (Oct "
        "24, 2024); Clifford Chance & MoFo case notes. "
        "Teaching beat (chronological, per Nico 2026-08-20): the "
        "previous slide ends on 'market definition would turn out to be "
        "crucial'; here reveal the ladder, then the share cards, then "
        "the internal quote, and close with the court decision as the "
        "punchline."))
    return slide


# --------------------------------------------------------------------------
# EXPANDED 2a/2b — Kroger–Albertsons (market definition)
# --------------------------------------------------------------------------

def exp_kroger_case(prs, page_num):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, "Module 1 · Candidates · Market Definition")
    _draw_action_title(slide,
                       "Market Definition Mini-Case: Kroger–Albertsons")
    box = _add_hierarchical_bullets(
        slide, left=MARGIN, top=Inches(1.50), width=Inches(7.9),
        height=Inches(5.3),
        items=[
            ("Oct 2022: Kroger agrees to buy Albertsons for $24.6B — the "
             "largest US supermarket merger ever proposed", 0),
            ([("Once again, everything would hinge on ", {}),
              ("how you define the market", {'bold': True}),
              (":", {})], 0, {}),
            ([("The firms: we are small next to ", {}),
              ("Walmart, Amazon, Costco", {'color': RED})], 1, {}),
            ([("FTC: the market is ", {}),
              ("“supermarkets”", {'color': RED}),
              (" — assessed city by city, where #1 was buying #2",
               {})], 1, {}),
        ],
        size=24, sub_size=22, line_spacing_pts=16,
        sub_line_spacing_pts=8)
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    _add_media_image(slide, "web_kroger.jpg",
                     left=Inches(8.75), top=Inches(1.70),
                     width=Inches(3.27), rounded=True)     # h ≈ 2.45
    _add_media_image(slide, "web_albertsons.jpg",
                     left=Inches(8.26), top=Inches(4.30),
                     width=Inches(4.25), rounded=True)     # h ≈ 2.43
    _photo_caption(slide, Inches(8.26), Inches(6.80), Inches(4.25))
    _footer(slide, page_num)
    _set_notes(slide, (
        "Expanded from the candidate list (verified 2026-08-20). Judge "
        "Adrienne Nelson (D. Or.) granted the FTC's preliminary "
        "injunction 10 Dec 2024 after a three-week hearing; a Washington "
        "state court blocked the deal under state law the same day. The "
        "proposed divestiture of 579 stores to C&S Wholesale was rejected "
        "as inadequate. Aftermath: Albertsons sued Kroger for the $600M "
        "termination fee (ongoing 2026). Teaching beat (chronological): "
        "this slide sets up the case and ends on 'market definition "
        "would turn out to be crucial'; the next slide resolves it.\n"
        "Photos: Wikimedia Commons (Kroger Marketplace, Athens GA; "
        "Albertsons, Dallas)."))
    return slide


def exp_costco_run(prs, page_num):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, "Module 1 · Candidates · Market Definition")
    _draw_action_title(slide,
                       "Is a “Costco Run” a Substitute for Everyday "
                       "Shopping?")
    # in-the-market box
    _add_rounded_filled_box(
        slide, Inches(0.75), Inches(1.62), Inches(6.9), Inches(0.95),
        "THE MARKET:  “supermarkets” — Kroger · Albertsons · Safeway …",
        fill=NAVY, text_color=WHITE, size=18)
    _add_text(slide, Inches(0.75), Inches(2.76), Inches(6.9), Inches(0.34),
              "Outside the market:", size=18, bold=True, color=GRAY,
              font="Calibri")
    outs = ["Club stores — Costco, Sam's Club",
            "Limited assortment — Aldi, Trader Joe's",
            "Dollar & convenience stores",
            "Online-only sellers"]
    for i, t in enumerate(outs):
        r, c = divmod(i, 2)
        _add_rounded_filled_box(
            slide, Inches(0.75 + c * 3.55), Inches(3.18 + r * 0.92),
            Inches(3.35), Inches(0.78), t,
            fill=FADED, text_color=WHITE, size=18, bold=False)
    # Costco photo
    _add_media_image(slide, "web_costco.jpg",
                     left=Inches(8.05), top=Inches(1.85),
                     width=Inches(4.55), rounded=True)     # h ≈ 2.57
    _photo_caption(slide, Inches(8.05), Inches(4.50), Inches(4.55),
                   "Photo: Wikimedia Commons")
    _quote_box(slide, MARGIN + Inches(0.35), Inches(5.15),
               RULE_W - Inches(0.7), Inches(1.15),
               "“A monthly trip to Costco to stock up … does not make a "
               "‘Costco run’ a reasonable substitute for a weekly one-stop "
               "visit to a supermarket.”",
               "Judge Adrienne Nelson, D. Or. (Dec 2024)")
    # main takeaway: gold bar (sources live in the notes)
    _add_rounded_filled_box(
        slide, Inches(0.90), Inches(6.44), Inches(11.53), Inches(0.60),
        "Dec 2024: federal and state courts block the deal — some "
        "substitution is not enough substitution",
        fill=GOLD, text_color=NAVY, size=19, bold=True, corner_pct=0.18)
    _footer(slide, page_num)
    _set_notes(slide, (
        "The court also recognized a broader 'large-format store' market "
        "that DOES include club stores — the narrow supermarkets market "
        "is where the presumption of illegality arose. Open point: "
        "whether Walmart supercenters sat inside the narrow market is "
        "not pinned down — keep supercenters off the slide. Teaching "
        "beat (chronological): market diagram, then the Nelson quote, "
        "then the Dec 2024 gold decision bar as the final click. "
        "On-slide source line dropped for space — cite: D. Or. "
        "preliminary injunction (Dec 10, 2024); SMU post-mortem; "
        "Grocery Dive."))
    return slide


# --------------------------------------------------------------------------
# EXPANDED 3a/3b — Netflix + Warner Bros (market definition)
# --------------------------------------------------------------------------

def exp_netflix_saga(prs, page_num):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, "Module 1 · Candidates · Market Definition")
    _draw_action_title(slide,
                       "Mini-Case: Netflix Bids for Warner Bros (2025–26)")
    box = _add_hierarchical_bullets(
        slide, left=MARGIN, top=Inches(1.48), width=RULE_W,
        height=Inches(0.6),
        items=[("Every side defined the market to fit its answer — "
                "Netflix, the DOJ, and 12 state AGs", 0)],
        size=22, line_spacing_pts=0)
    box.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    _add_media_image(slide, "ic_s20_rId2.png",
                     left=Inches(2.35), top=Inches(2.20),
                     width=Inches(3.7), rounded=True)      # h ≈ 2.55
    _add_media_image(slide, "web_wb.jpg",
                     left=Inches(7.35), top=Inches(2.20),
                     width=Inches(3.8), rounded=True)      # h ≈ 2.53
    _photo_caption(slide, Inches(7.35), Inches(4.78), Inches(3.8),
                   "Photo: Wikimedia Commons")
    # timeline
    line_y = Inches(5.65)
    _add_arrow(slide, (Inches(0.9), line_y), (Inches(12.6), line_y),
               color=NAVY, weight_pt=2.5, head=True)
    stops = [
        ("Dec 2025", "Netflix agrees to buy WB studios + HBO Max "
                     "(~$83B reported)"),
        ("Feb 2026", "DOJ opens an antitrust probe"),
        ("Feb 26, 2026", "Paramount outbids (~$110.9B); Netflix "
                         "withdraws"),
        ("Jul 2026", "12 state AGs sue Paramount–WBD — over theatrical "
                     "film + cable markets"),
    ]
    for i, (date, desc) in enumerate(stops):
        x = Inches(1.7 + i * 3.0)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     int(x - Inches(0.10)),
                                     int(line_y - Inches(0.10)),
                                     int(Inches(0.20)), int(Inches(0.20)))
        dot.fill.solid(); dot.fill.fore_color.rgb = GOLD
        dot.line.color.rgb = NAVY; dot.line.width = Pt(1.0)
        dot.shadow.inherit = False
        _add_text(slide, x - Inches(1.4), line_y - Inches(0.52),
                  Inches(2.8), Inches(0.34), date, size=18, bold=True,
                  color=NAVY, font="Calibri", align=PP_ALIGN.CENTER)
        _add_text(slide, x - Inches(1.45), line_y + Inches(0.18),
                  Inches(2.9), Inches(1.15), desc, size=16, color=GRAY,
                  font="Calibri", align=PP_ALIGN.CENTER)
    _footer(slide, page_num)
    _set_notes(slide, (
        "Timeline verified 2026-08-20: WBD–Netflix merger agreement 4 "
        "Dec 2025; DOJ civil investigative demands Feb 2026 ('may "
        "substantially lessen competition'); WBD board declared "
        "Paramount Skydance's ~$110.9B offer superior 26 Feb 2026 and "
        "Netflix withdrew; 12 state AGs (led by CA's Bonta) sued 13 Jul "
        "2026 in N.D. Cal. — alleging harm in wide-release theatrical "
        "distribution (~27% combined), top-grossing theatrical (>30%), "
        "and basic-cable licensing (~27%), not in streaming.\n"
        "Photo: Warner Bros. Studios water tower, Wikimedia Commons."))
    return slide


def exp_netflix_chart(prs, page_num):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, "Module 1 · Candidates · Market Definition")
    _draw_action_title(slide,
                       "9% of TV Time — or 33% of Streaming? Pick Your "
                       "Market")
    # Nielsen share-of-TV-time bars (Netflix investor deck, Oct 2025)
    data = [("YouTube", 12.9, NAVY), ("Disney", 11.4, NAVY),
            ("NBCU", 8.6, NAVY), ("Fox", 8.4, NAVY),
            ("Paramount", 8.2, NAVY), ("Netflix", 8.0, GOLD),
            ("WBD", 5.6, NAVY)]
    base_y = 5.55
    scale = 0.27          # inches per share point
    x0, bw, gap = 1.05, 0.92, 0.33
    hbo = 1.2             # HBO Max share inside WBD
    for i, (name, val, fill) in enumerate(data):
        x = x0 + i * (bw + gap)
        h = val * scale
        if name == "WBD":
            _add_rect(slide, Inches(x), Inches(base_y - h + hbo * scale),
                      Inches(bw), Inches(h - hbo * scale), NAVY)
            _add_rect(slide, Inches(x), Inches(base_y - h),
                      Inches(bw), Inches(hbo * scale), GOLD)
            _add_text(slide, Inches(x - 0.55), Inches(base_y + 0.40),
                      Inches(2.0), Inches(0.26), "(HBO Max: 1.2)",
                      size=12, bold=True, color=GOLD, font="Calibri",
                      align=PP_ALIGN.CENTER)
        else:
            _add_rect(slide, Inches(x), Inches(base_y - h),
                      Inches(bw), Inches(h), fill)
        _add_text(slide, Inches(x - 0.25), Inches(base_y - h - 0.34),
                  Inches(bw + 0.5), Inches(0.30),
                  ("%.1f" % val) if name != "WBD" else "5.6",
                  size=16, bold=True,
                  color=GOLD if name == "Netflix" else NAVY,
                  font="Calibri", align=PP_ALIGN.CENTER)
        _add_text(slide, Inches(x - 0.30), Inches(base_y + 0.08),
                  Inches(bw + 0.6), Inches(0.3), name, size=14,
                  bold=(name in ("Netflix", "WBD")), color=NAVY,
                  font="Calibri", align=PP_ALIGN.CENTER)
    _add_rect(slide, Inches(x0 - 0.15), Inches(base_y),
              Inches(7 * (bw + gap) + 0.1), Inches(0.022), NAVY)
    _add_text(slide, Inches(x0 - 0.15), Inches(6.12),
              Inches(7 * (bw + gap)), Inches(0.3),
              "Share of US TV time, Oct 2025 (Nielsen, via Netflix's "
              "investor deck)", size=13, italic=True, color=GRAY,
              font="Calibri", align=PP_ALIGN.CENTER)
    # Netflix's framing + the rival framing
    _add_convention_box(
        slide, Inches(9.55), Inches(1.60), Inches(3.15), Inches(2.05),
        runs=[("Netflix's market: all TV time", {'bold': True, 'size': 18}),
              ("Netflix + HBO Max = 9.2% — “#3 behind YouTube and Disney”",
               {'size': 18, 'newline': True})])
    box2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, int(Inches(9.55)), int(Inches(3.90)),
        int(Inches(3.15)), int(Inches(2.05)))
    try:
        box2.adjustments[0] = 0.10
    except Exception:
        pass
    box2.fill.solid(); box2.fill.fore_color.rgb = WHITE
    box2.line.color.rgb = RED; box2.line.width = Pt(1.5)
    box2.shadow.inherit = False
    _add_drop_shadow(box2)
    tb = slide.shapes.add_textbox(int(Inches(9.72)), int(Inches(4.00)),
                                  int(Inches(2.85)), int(Inches(1.85)))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = "Narrow SVOD market: "
    r1.font.name = "Calibri"; r1.font.size = Pt(18)
    r1.font.bold = True; r1.font.color.rgb = RED
    r2 = p.add_run()
    r2.text = "~33% combined — presumptively illegal concentration"
    r2.font.name = "Calibri"; r2.font.size = Pt(18)
    r2.font.color.rgb = NAVY
    _add_rounded_filled_box(
        slide, Inches(0.90), Inches(6.44), Inches(11.53), Inches(0.60),
        "Both numbers are correct arithmetic — the market definition "
        "does all the work",
        fill=GOLD, text_color=NAVY, size=19, bold=True, corner_pct=0.18)
    _footer(slide, page_num)
    _set_notes(slide, (
        "On-slide source line dropped for space — cite: Netflix "
        "investor deck (Dec 2025, Nielsen Oct 2025); ProMarket / "
        "Wolfram (Feb 2026). "
        "Bars: Nielsen 'Share of US TV Time by Distributor', Oct 2025, "
        "as published in Netflix's own investor deck (Dec 2025): YouTube "
        "12.9, Disney 11.4, NBCU 8.6, Fox 8.4, Paramount 8.2, Netflix "
        "8.0, WBD 5.6 (of which HBO/HBO Max 1.2) — Netflix + HBO Max = "
        "9.2%. SVOD framing: Wolfram/ProMarket (5 Feb 2026): combined "
        "~33% of SVOD, HHI 2,055, delta 829 — over the structural "
        "presumption. Build: bars first, Netflix framing box, then the "
        "red SVOD box as the punchline."))
    return slide


# --------------------------------------------------------------------------
# EXPANDED 4a/4b — United and the marginal flight (MB = MC)
# --------------------------------------------------------------------------

def exp_united_case(prs, page_num):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide,
                     "Module 1 · Candidates · Cost-Benefit and Marginal "
                     "Analysis")
    _draw_action_title(slide,
                       "Mini-Case: Airlines and the 2026 Fuel Shock")
    box = _add_hierarchical_bullets(
        slide, left=MARGIN, top=Inches(1.55), width=Inches(6.8),
        height=Inches(3.8),
        items=[
            ("Mar 2026: jet fuel “more than doubled in the last three "
             "weeks” (Kirby memo)", 0),
            ("United cuts ~5% of planned capacity:", 0),
            ("~3pp off-peak, midweek and red-eye flying", 1),
            ("~1pp Chicago O'Hare; ~1pp Tel Aviv / Dubai", 1),
            ([("Delta follows: off-peak and red-eyes are ", {}),
              ("“15% to 20% less valuable on a net revenue basis”",
               {'color': RED}),
              (" (CCO Esposito)", {})], 0, {}),
            ("Aircraft orders and 2027 capacity: unchanged", 0),
        ],
        size=21, sub_size=19, line_spacing_pts=10)
    box.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    _add_media_image(slide, "web_united.jpg",
                     left=Inches(7.65), top=Inches(2.05),
                     width=Inches(5.0), rounded=True)      # h ≈ 2.43
    _photo_caption(slide, Inches(7.65), Inches(4.55), Inches(5.0),
                   "Photo: Wikimedia Commons")
    _quote_box(slide, MARGIN + Inches(0.35), Inches(5.55),
               RULE_W - Inches(0.7), Inches(1.10),
               "“There's no point in burning cash in the near term on "
               "flying that just can't absorb these fuel costs. … Nothing "
               "changes about our longer-term plans.”",
               "Scott Kirby, United CEO (staff memo, Mar 2026)")
    _source_line(slide, "Kirby memo via Fox Business / CNBC (Mar 2026); "
                        "Travel Market Report (Apr 2026)")
    _footer(slide, page_num)
    _set_notes(slide, (
        "Verified 2026-08-20 from coverage of the 20 Mar 2026 memo: ~5% "
        "capacity cut; United modeling oil to $175/bbl and >$100 through "
        "2027 (≈$11B added annual fuel cost — 'more than twice the "
        "profit it earned in its best year ever'). Delta (Apr 2026, EVP "
        "Joe Esposito) targeted off-peak/edge-of-day/red-eyes. Short "
        "run vs. long run stated by the CEO himself.\n"
        "Photo: United 787-8, Wikimedia Commons."))
    return slide


def exp_united_chart(prs, page_num):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide,
                     "Module 1 · Candidates · Cost-Benefit and Marginal "
                     "Analysis")
    _draw_action_title(slide,
                       "When MC Jumps, Cut the Marginal Flight — Not the "
                       "Fleet")
    fig = SimpleFig(3.3, 5.95, 6.6, 3.9, 10, 10)
    _fig_axes(slide, fig, y_title="$ per flight", x_title="Flights",
              label_size=16)
    # MR of flights, ranked: y = 8.9167 - 0.8333x
    _fig_line(slide, fig, (0.5, 8.5), (9.5, 1.0), color=GREEN_MB,
              weight_pt=2.75)
    _add_text(slide, fig.x(1.6), fig.y(7.35), Inches(2.9), Inches(0.32),
              "MR of flights (ranked)", size=16, bold=True, italic=True,
              color=GREEN_MB, font="Calibri")
    # MC before / after (labels on the empty left side of the plot)
    _fig_line(slide, fig, (0, 3.2), (9.7, 3.2), color=NAVY,
              weight_pt=2.25, dash='dash')
    _add_text(slide, fig.x(0.3), fig.y(2.95), Inches(2.7), Inches(0.32),
              "MC — normal fuel", size=16, bold=True, italic=True,
              color=NAVY, font="Calibri")
    _fig_line(slide, fig, (0, 4.6), (9.7, 4.6), color=RED, weight_pt=2.5)
    _add_text(slide, fig.x(0.3), fig.y(5.45), Inches(2.7), Inches(0.32),
              "MC — fuel ×2", size=16, bold=True, italic=True,
              color=RED, font="Calibri")
    # crossings: x0 = 6.86 (old), x1 = 5.18 (new)
    _fig_guide(slide, fig, (6.86, 3.2), to_y=False, color=GRAY)
    _fig_guide(slide, fig, (5.18, 4.6), to_y=False, color=GRAY)
    _fig_xlab(slide, fig, 6.86, "Q0*", size=16)
    _fig_xlab(slide, fig, 5.18, "Q1*", size=16)
    # the cut band (inside the plot, clear of the axis labels)
    arr = _add_arrow(slide, (fig.x(5.18), fig.y(0.85)),
                     (fig.x(6.86), fig.y(0.85)),
                     color=GOLD, weight_pt=3.0, head=True)
    try:
        from lxml import etree as _ET
        from pptx.oxml.ns import qn as _qn
        ln = arr.line._get_or_add_ln()
        hd = _ET.SubElement(ln, _qn('a:headEnd'))
        hd.set('type', 'triangle'); hd.set('w', 'med'); hd.set('h', 'med')
    except Exception:
        pass
    _add_text(slide, Inches(5.52), fig.y(2.60), Inches(3.5), Inches(0.66),
              "the ~5% cut:\noff-peak and red-eyes first", size=15,
              bold=True, color=GOLD, font="Calibri",
              align=PP_ALIGN.CENTER)
    _add_convention_box(
        slide, MARGIN + Inches(0.35), Inches(6.38),
        RULE_W - Inches(0.7), Inches(0.70),
        runs=[("Shed exactly the flights between the old and new MB = MC "
               "crossing — short run only; the fleet (long run) stays.",
               {'size': 18})],
        align=PP_ALIGN.CENTER)
    _footer(slide, page_num)
    _set_notes(slide, (
        "Schematic (values illustrative): flights ranked by marginal "
        "revenue give a downward MR line; fuel is the textbook marginal "
        "cost, so a fuel-price doubling shifts MC up and the crossing "
        "left. The flights between Q0* and Q1* — the red-eyes and "
        "off-peak departures 15–20% below average revenue — are exactly "
        "the ones cancelled. Lands right after the exercise-hours and "
        "continuous MB=MC slides (57–59)."))
    return slide


# --------------------------------------------------------------------------
# EXPANDED 5a/5b/5c — DRAM / AI memory (supply and demand, 3 slides)
# --------------------------------------------------------------------------

def exp_dram_case(prs, page_num):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, "Module 1 · Candidates · Supply and Demand")
    _draw_action_title(slide,
                       "Mini-Case: AI Eats the Memory Supply (2025–26)")
    box = _add_hierarchical_bullets(
        slide, left=MARGIN, top=Inches(1.50), width=Inches(7.3),
        height=Inches(5.3),
        items=[
            ("2025: the AI datacenter buildout explodes demand for "
             "high-bandwidth memory (HBM)", 0),
            ("New chip fabs take years — industry capacity is "
             "essentially fixed for 2–3 years", 0),
            ([("Oct 2025: SK Hynix — DRAM, NAND and HBM capacity ", {}),
              ("“essentially sold out” for 2026", {'color': RED})],
             0, {}),
            ([("The same wafers can become AI memory or ordinary laptop "
               "memory", {'bold': True}),
              (" — that allocation would drive the story", {})], 0, {}),
        ],
        size=24, sub_size=22, line_spacing_pts=18)
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    _add_media_image(slide, "web_datacenter.jpg",
                     left=Inches(7.95), top=Inches(2.30),
                     width=Inches(4.6), rounded=True)      # h ≈ 3.07
    _photo_caption(slide, Inches(7.95), Inches(5.42), Inches(4.6),
                   "Photo: Wikimedia Commons")
    _footer(slide, page_num)
    _set_notes(slide, (
        "Setup slide (chronology-first): builds the situation and ends "
        "on the wafer-allocation flag; the S/D analysis and the price "
        "outcome follow on the next two slides. Confirmed: SK Hynix's "
        "'essentially sold out' statement is from its Oct 2025 earnings "
        "call, corroborated across independent outlets.\n"
        "Sources:\n"
        "https://www.networkworld.com/article/4113772/samsung-warns-of-"
        "memory-shortages-driving-industry-wide-price-surge-in-2026.html\n"
        "https://www.trendforce.com/presscenter/news/20260105-12860.html"))
    return slide


def exp_dram_sd(prs, page_num):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, "Module 1 · Candidates · Supply and Demand")
    _draw_action_title(slide, "Two Linked Markets, One Pool of Wafers")

    _add_rounded_filled_box(slide, Inches(0.85), Inches(1.52),
                            Inches(4.9), Inches(0.52),
                            "HBM — AI memory",
                            fill=NAVY, text_color=WHITE, size=18)
    _add_rounded_filled_box(slide, Inches(7.45), Inches(1.52),
                            Inches(4.9), Inches(0.52),
                            "Consumer DRAM (DDR5)",
                            fill=NAVY, text_color=WHITE, size=18)

    # Panel 1 — HBM: demand explodes against a steep (near-fixed) supply.
    fig1 = SimpleFig(1.55, 6.00, 3.7, 3.35, 10, 10)
    _fig_axes(slide, fig1, label_size=13)
    _fig_guide(slide, fig1, (4.44, 2.56), color=GRAY)
    _fig_guide(slide, fig1, (5.44, 6.56), color=GRAY)
    _fig_line(slide, fig1, (4.0, 0.8), (6.0, 8.8), color=STEEL,
              weight_pt=2.75)
    _fig_curve_label(slide, fig1, 6.15, 9.2, "S", color=NAVY)
    _fig_line(slide, fig1, (0.6, 6.4), (6.4, 0.6), color=GOLD,
              weight_pt=2.5)
    _fig_curve_label(slide, fig1, 6.55, 1.0, "D", color=GOLD)
    _fig_line(slide, fig1, (3.6, 8.4), (9.6, 2.4), color=GREEN_BR,
              weight_pt=2.5, dash='dash')
    _fig_curve_label(slide, fig1, 9.15, 3.1, "D’", color=GREEN_BR)
    _add_arrow(slide, (fig1.x(3.4), fig1.y(4.4)),
               (fig1.x(4.9), fig1.y(5.6)),
               color=GREEN_BR, weight_pt=2.0, head=True)
    _fig_ylab(slide, fig1, 2.56, "P0", size=14)
    _fig_ylab(slide, fig1, 6.56, "P1", size=14)

    # Panel 2 — consumer DRAM: supply shifts LEFT as wafers leave.
    fig2 = SimpleFig(8.15, 6.00, 3.7, 3.35, 10, 10)
    _fig_axes(slide, fig2, label_size=13)
    _fig_guide(slide, fig2, (4.5, 3.5), color=GRAY)
    _fig_guide(slide, fig2, (3.0, 5.0), color=GRAY)
    _fig_line(slide, fig2, (1.5, 0.5), (8.5, 7.5), color=STEEL,
              weight_pt=2.5)
    _fig_curve_label(slide, fig2, 8.65, 7.9, "S", color=NAVY)
    _fig_line(slide, fig2, (0.8, 2.8), (6.8, 8.8), color=BLUE_PED,
              weight_pt=2.5, dash='dash')
    _fig_curve_label(slide, fig2, 6.0, 9.2, "S’", color=BLUE_PED)
    _fig_line(slide, fig2, (0.8, 7.2), (7.4, 0.6), color=GOLD,
              weight_pt=2.5)
    _fig_curve_label(slide, fig2, 7.55, 1.0, "D", color=GOLD)
    _add_arrow(slide, (fig2.x(6.2), fig2.y(5.6)),
               (fig2.x(4.8), fig2.y(6.9)),
               color=BLUE_PED, weight_pt=2.0, head=True)
    _fig_ylab(slide, fig2, 3.5, "P0", size=14)
    _fig_ylab(slide, fig2, 5.0, "P1", size=14)
    _fig_xlab(slide, fig2, 4.5, "Q0", size=14)
    _fig_xlab(slide, fig2, 3.0, "Q1", size=14)

    # the reallocation arrow between the panels
    _add_arrow_shape(slide, Inches(5.80), Inches(3.55), Inches(1.55),
                     Inches(0.55), direction="right", fill=GOLD)
    _add_text(slide, Inches(5.48), Inches(4.20), Inches(2.20),
              Inches(1.05), "wafers reallocated to the higher-margin use",
              size=15, bold=True, color=GOLD, font="Calibri",
              align=PP_ALIGN.CENTER)

    _add_convention_box(
        slide, MARGIN + Inches(0.35), Inches(6.40),
        RULE_W - Inches(0.7), Inches(0.68),
        runs=[("Nothing physically broke — the marginal wafer flowed to "
               "its best-paying use: opportunity cost as a supply-side "
               "force.", {'size': 18})],
        align=PP_ALIGN.CENTER)
    _footer(slide, page_num)
    _set_notes(slide, (
        "Two-panel S/D schematic (values illustrative). Left: AI demand "
        "shifts right hard against a steep, near-fixed short-run supply "
        "— price does almost all the adjusting. Right: the SAME wafer "
        "capacity is reallocated toward high-margin HBM, so the "
        "consumer-DRAM supply curve shifts LEFT with demand unchanged — "
        "price up, quantity down. A demand shock in one market becomes "
        "a supply shock in the linked one. Teaching beat: left panel "
        "first (D→D'), then the gold reallocation arrow, then the right "
        "panel (S→S'), takeaway box last."))
    return slide


def exp_dram_outcome(prs, page_num):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, "Module 1 · Candidates · Supply and Demand")
    _draw_action_title(slide,
                       "Resolution: Prices Double — High Prices Can't "
                       "Fix This Quickly")
    box = _add_hierarchical_bullets(
        slide, left=MARGIN, top=Inches(1.50), width=Inches(7.5),
        height=Inches(4.0),
        items=[
            ("Q1 2026: conventional DRAM contract prices roughly double "
             "in a single quarter (TrendForce, reported)", 0),
            ("Samsung reportedly reprices a 32GB DDR5 module $149 → "
             "$239 in one step (+60%)", 0),
            ("HBM's share of DRAM wafer output: ~19% → ~23% (2026, "
             "reported); makers signal tightness into 2028", 0),
            ([("New fabs take 2–3 years — quantity can't respond, so "
               "price does all the work", {'bold': True})], 0, {}),
        ],
        size=22, sub_size=20, line_spacing_pts=14)
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    _add_media_image(slide, "web_h100.jpg",
                     left=Inches(8.30), top=Inches(1.70),
                     width=Inches(4.35), rounded=True)     # h ≈ 2.45
    _add_media_image(slide, "web_dram.jpg",
                     left=Inches(8.95), top=Inches(4.28),
                     width=Inches(3.05), rounded=True)     # h ≈ 2.44
    _photo_caption(slide, Inches(8.30), Inches(6.78), Inches(4.35))
    disc = slide.shapes.add_textbox(int(MARGIN), int(Inches(5.72)),
                                    int(Inches(7.6)), int(Inches(0.95)))
    tf = disc.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = "Discussion:  "
    r1.font.name = "Calibri"; r1.font.size = Pt(18)
    r1.font.bold = True; r1.font.color.rgb = GOLD
    r2 = p.add_run()
    r2.text = ("No factory burned down, yet laptop-memory supply "
               "contracted. What contract terms protect you when your "
               "input has a better-paying alternative use?")
    r2.font.name = "Calibri"; r2.font.size = Pt(18)
    r2.font.italic = True; r2.font.color.rgb = NAVY
    _add_text(slide, MARGIN, Inches(6.78), Inches(7.6), Inches(0.3),
              "Sources: TrendForce press center (2026, reported); "
              "Network World; SK Hynix earnings call (Oct 2025)",
              size=11, italic=True, color=GRAY, font="Calibri")
    _footer(slide, page_num)
    _set_notes(slide, (
        "Resolution slide (final beat when animated). VERIFICATION: the "
        "direction and mechanism are corroborated across independent "
        "outlets, but the MAGNITUDES (the ~93–98% QoQ Q1-2026 contract-"
        "price move, the $149→$239 Samsung repricing, the 19%→23% HBM "
        "wafer share) are press-reported TrendForce figures surfaced via "
        "search, not primary-fetched — re-verify the TrendForce releases "
        "before this graduates into the main deck.\n"
        "Photos: NVIDIA H100 accelerators (where the wafers went) and SK "
        "Hynix DDR5 modules — Wikimedia Commons.\n"
        "Sources:\n"
        "https://www.trendforce.com/presscenter/news/20260601-13070.html\n"
        "https://www.trendforce.com/presscenter/news/20260105-12860.html\n"
        "https://www.networkworld.com/article/4113772/samsung-warns-of-"
        "memory-shortages-driving-industry-wide-price-surge-in-2026.html"))
    return slide


# --------------------------------------------------------------------------
# CANDIDATES — curated from the verified research results (2026-08-20).
# Each: (concept, title, facts, angle, discussion, sources, extras-dict).
# Facts marked "reported" are press-reported, not primary-confirmed;
# full source URLs + verification flags live in the speaker notes.
# --------------------------------------------------------------------------
CANDIDATES = [

    # ---- Market definition -------------------------------------------------
    ("Market Definition",
     "“Accessible Luxury”: The Market Definition That Killed an $8.5B Deal",
     [
         "Tapestry (Coach, Kate Spade) agreed to buy Capri (Michael Kors, "
         "Versace) for $8.5B (Aug 2023)",
         "FTC's market: “accessible luxury” handbags — roughly $100 to "
         "under $1,000",
         "Combined share 59% on FTC data — up to 83% on Tapestry's own "
         "data; trivial in “all handbags”",
         "Court blocked the deal (Oct 2024) citing the firms' own words: "
         "“Nobody says ‘should I buy a LV bag or a Coach bag?’”",
         "Deal abandoned Nov 2024",
     ],
     "The market boundary, not the conduct, decided the case — and the "
     "firms' own strategy documents drew the boundary. Natural companion "
     "to the ADM mini-case (slide 19).",
     "If the FTC subpoenaed your strategy deck tomorrow, would your own "
     "market definition help you or convict you?",
     "Judge Rochon opinion, S.D.N.Y. (Oct 24, 2024); Clifford Chance & "
     "MoFo case notes; Capri 10-K (termination)",
     {"visual": "Coach / Michael Kors product photos + a price-tier "
                "ladder graphic (mass — accessible luxury — true luxury)",
      "notes": (
        "All figures verified against the cited case notes on 2026-08-20. "
        "Shares: FTC economist 58.7% (third-party data); Capri internal "
        "documents implied 77%, Tapestry internal data 83%. Preliminary "
        "injunction granted 24 Oct 2024 (169-page opinion); merger "
        "agreement terminated 13 Nov 2024.\n"
        "Sources:\n"
        "https://www.cliffordchance.com/content/dam/cliffordchance/"
        "briefings/2024/10/antitrust-has-come-into-fashion-us-federal-"
        "trade-commission-wins-preliminary-injunction-against-tapestrys-"
        "acquisition-of-capri.pdf\n"
        "https://www.mofo.com/resources/insights/241029-tapestry-capri-"
        "handbag-merger-halted-by-s-d-n-y\n"
        "https://www.thefashionlaw.com/judge-preliminarily-blocks-"
        "tapestry-capri-merger-in-win-for-the-ftc/")}),

    ("Market Definition",
     "Is a “Costco Run” a Substitute? The $24.6B Supermarket Merger",
     [
         "Kroger–Albertsons ($24.6B): largest proposed US supermarket "
         "merger; blocked Dec 2024",
         "Court adopted a narrow “supermarkets” market — excluding "
         "Costco, Aldi, Trader Joe's, dollar stores",
         "Judge Nelson: a monthly “Costco run” is no substitute for the "
         "weekly one-stop supermarket trip",
         "Markets assessed city by city, not nationally; deal abandoned "
         "Dec 2024",
     ],
     "The slide-18 price test in plain English: SOME substitution is not "
     "ENOUGH substitution. Also introduces the geographic dimension of "
     "market definition.",
     "What evidence would prove Costco does discipline supermarket "
     "prices — scanner data, a Costco-opening natural experiment, "
     "diversion surveys?",
     "D. Or. preliminary injunction (Dec 10, 2024); Grocery Dive; "
     "Washington State Standard; SMU post-mortem",
     {"visual": "Kroger vs. Costco storefront photos; city-dot map for "
                "the geographic markets",
      "notes": (
        "Verified 2026-08-20. Washington state court blocked the deal "
        "the same day under state law. Divestiture of 579 stores to C&S "
        "rejected as inadequate. Aftermath: Albertsons sued Kroger for "
        "the $600M termination fee (Delaware Chancery, ongoing 2026). "
        "Open point: whether Walmart supercenters sat inside the narrow "
        "market is NOT pinned down — keep supercenters off the slide.\n"
        "Sources:\n"
        "https://www.grocerydive.com/news/kroger-albertsons-merger-ftc-"
        "judge-preliminary-injunction/727132/\n"
        "https://washingtonstatestandard.com/2024/12/10/judges-in-oregon-"
        "washington-block-kroger-albertsons-supermarket-merger/\n"
        "https://www.smu.edu/-/media/site/cox/faculty/research/foxedward/"
        "lessons-learned-from-the-kroger-albertsons-merger-case.pdf")}),

    ("Market Definition",
     "Netflix + Warner Bros: One Deal, Three Market Definitions",
     [
         "Dec 2025: Netflix agreed to buy Warner Bros studios + HBO Max "
         "(reported ~$83B)",
         "Netflix's investor slide — “US TV time” market: Netflix + WB "
         "= 9.2%, behind YouTube (12.9%) and Disney (11.4%)",
         "Narrow SVOD market instead: ~33% combined — presumptively "
         "illegal concentration",
         "Feb 2026: DOJ probe opened; Paramount outbid at ~$110.9B; 12 "
         "state AGs sued — over theatrical film and cable markets, not "
         "streaming",
     ],
     "A live upgrade for the existing “Define Netflix's market” slides "
     "(20 and 68): every player defined the market to fit its answer, "
     "and Netflix's own 9.2% chart is a ready-made teaching artifact.",
     "9.2% of TV time and 33% of SVOD are both correct arithmetic — "
     "which better predicts pricing power?",
     "Netflix investor deck (Dec 2025); ProMarket (Feb 2026); CA AG "
     "complaint (Jul 2026)",
     {"visual": "Native rebuild of Netflix's “Share of US TV Time” bar "
                "chart (public investor PDF)",
      "notes": (
        "Verified 2026-08-20. Timeline: WBD–Netflix agreement 4 Dec "
        "2025; Nielsen Oct-2025 shares as on Netflix's slide; DOJ civil "
        "investigative demands Feb 2026; WBD board declared Paramount "
        "Skydance's ~$110.9B offer superior 26 Feb 2026, Netflix "
        "withdrew; 12 state AGs (led by CA's Bonta) sued 13 Jul 2026 in "
        "N.D. Cal. over wide-release theatrical distribution (~27% "
        "combined), top-grossing theatrical (>30%), and basic-cable "
        "licensing (~27%). SVOD estimate: Wolfram/ProMarket — combined "
        "~33%, HHI 2,055, delta 829.\n"
        "Sources:\n"
        "https://s22.q4cdn.com/959853165/files/doc_events/2025/Dec/08/"
        "NFLX-WB-Share.pdf\n"
        "https://www.promarket.org/2026/02/05/netflix-appears-to-face-"
        "greater-antitrust-barriers-to-acquiring-warner-bros-discovery-"
        "than-paramount/\n"
        "https://oag.ca.gov/news/press-releases/attorney-general-bonta-"
        "files-lawsuit-block-110-billion-warner-brosparamount")}),

    # ---- Supply and demand -------------------------------------------------
    ("Supply and Demand",
     "Eggs 2024–2026: A Supply Shock in Fast-Forward",
     [
         "Avian flu culled the laying flock: supply shifts left, demand "
         "unchanged",
         "Retail eggs: $2.52/dozen (Jan 2024) → $6.23 (Mar 2025) → $2.19 "
         "(Jul 2026) — BLS data",
         "Inelastic demand: a supply loss well under 15% produced a 147% "
         "price spike",
         "Rationing without prices: Trader Joe's 1-dozen limits; Waffle "
         "House's transparent 50¢-per-egg surcharge (Feb–Jul 2025)",
     ],
     "The cleanest fully-verified S/D episode available — and hens "
     "regrow in months, so the whole cycle reversed inside two years. "
     "Slots naturally after the tea and avocado mini-cases.",
     "Supply fell under 15%, price rose 147% — what does that ratio say "
     "about the demand curve, and would you pass the cost through?",
     "BLS series APU0000708111 (retrieved 2026-08-20); USDA (2025); "
     "GMA (purchase limits); CNN / Fortune (Waffle House)",
     {"visual": "Native line chart of the BLS retail egg-price series "
                "2024–26 (data already pulled)",
      "notes": (
        "Price series pulled directly from the BLS public API on "
        "2026-08-20 (series APU0000708111, Grade A large, US city "
        "average) — fully verified. Secondary flavor: >150M birds culled "
        "by early Feb 2025; purchase limits at Trader Joe's (1 dozen), "
        "Sprouts (4), Costco (5 online); Waffle House surcharge Feb 4 – "
        "early Jul 2025. Optional dark twist for class: DOJ + 17 states' "
        "Jun 2026 civil settlement with Cal-Maine et al. over Urner "
        "Barry benchmark manipulation ($3.3M + 53M donated eggs, no "
        "admission) — separates scarcity pricing from manipulation.\n"
        "Sources:\n"
        "https://api.bls.gov/publicAPI/v1/timeseries/data/APU0000708111\n"
        "https://www.usda.gov/about-usda/news/press-releases/2025/03/20/"
        "usda-update-progress-five-pronged-strategy-combat-avian-flu-and-"
        "lower-egg-prices\n"
        "https://www.goodmorningamerica.com/food/story/trader-joes-"
        "costco-sprouts-limit-purchases-fresh-eggs-118654037\n"
        "https://www.cnn.com/2025/02/04/food/waffle-house-egg-surcharge/"
        "index.html")}),

    ("Supply and Demand",
     "AI Ate the Memory Supply: DRAM Prices 2025–2026",
     [
         "AI datacenter demand exploded; chip-fab capacity is fixed for "
         "2–3 years",
         "DRAM contract prices reportedly ~doubled in Q1 2026 "
         "(TrendForce); SK Hynix: 2026 capacity “essentially sold out”",
         "Nothing physically broke: high-margin AI memory (HBM) outbid "
         "consumer memory for the SAME wafers",
         "A demand explosion in one market = a supply contraction in "
         "the linked one",
     ],
     "Opportunity cost as a supply-side force — capacity flowed to the "
     "better-paying use. Bridges the S/D block to the opportunity-cost "
     "block. Magnitudes are press-reported: verify TrendForce releases "
     "before adoption.",
     "No factory burned down, yet laptop-memory supply contracted. What "
     "contract terms protect you when your input has a better-paying "
     "alternative use?",
     "TrendForce press center (2026, reported); Network World (Samsung "
     "warning); SK Hynix earnings call (Oct 2025)",
     {"visual": "Two-panel native S/D diagram: AI demand shifts right → "
                "consumer-DRAM supply shifts left",
      "notes": (
        "Direction and mechanism corroborated across independent "
        "outlets; the specific percentage moves (~93–98% QoQ Q1 2026, "
        "~58–63% Q2) are TrendForce figures surfaced via search, NOT "
        "primary-fetched — verify before putting magnitudes on a slide. "
        "HBM reportedly ~23% of DRAM wafer output in 2026 vs ~19% prior "
        "year. SK Hynix 'essentially sold out' from its Oct 2025 "
        "earnings call.\n"
        "Sources:\n"
        "https://www.trendforce.com/presscenter/news/20260601-13070.html\n"
        "https://www.trendforce.com/presscenter/news/20260105-12860.html\n"
        "https://www.networkworld.com/article/4113772/samsung-warns-of-"
        "memory-shortages-driving-industry-wide-price-surge-in-2026.html")}),

    # ---- Opportunity costs -------------------------------------------------
    ("Opportunity Costs",
     "The $100M Question: What It Costs to Stay",
     [
         "Altman (Jun 2025): Meta offered OpenAI staff “$100 million "
         "signing bonuses” — confirmed quote",
         "Meta's CTO: multi-year stock packages, not sign-on bonuses — "
         "“the market's hot. It's not that hot.”",
         "NYT-reported: one 24-year-old researcher's Meta package "
         "~$250M over 4 years (reported, not confirmed)",
         "A researcher who declines $100M pays $100M — in implicit "
         "cost — to stay",
     ],
     "Retention is a purchase at the market price: the outside offer is "
     "the implicit cost of staying, whether or not payroll ever records "
     "it. Companion to the full-economic-cost-of-an-MBA discussion.",
     "Your star engineer declines an offer at 3× the current package. "
     "Has your cost of employing her/him changed — and what if the "
     "accountants say no?",
     "Bloomberg / CNBC (Jun 2025); TechCrunch (Jun 2025); NYT (Aug "
     "2025, reported)",
     {"visual": "Explicit-vs-implicit cost ledger + Meta / OpenAI logos",
      "notes": (
        "Confirmed: Altman's Uncapped-podcast statement (17 Jun 2025) "
        "and Bosworth's all-hands pushback (via The Verge/TechCrunch); "
        "researcher Lucas Beyer publicly denied receiving $100M sign-on. "
        "Reported only: the ~$250M/4-year package (~$100M year one) for "
        "Matt Deitke, NYT-attributed via relays; Meta total outlay >$1B "
        "on ~50 hires. Do not put the $250M on a main-deck slide without "
        "reading the NYT piece.\n"
        "Sources:\n"
        "https://www.bloomberg.com/news/articles/2025-06-17/altman-says-"
        "meta-offered-openai-staffers-100-million-bonuses\n"
        "https://www.cnbc.com/2025/06/18/sam-altman-says-meta-tried-to-"
        "poach-openai-staff-with-100-million-bonuses-mark-zuckerberg.html\n"
        "https://techcrunch.com/2025/06/27/meta-is-offering-multimillion-"
        "dollar-pay-for-ai-researchers-but-not-100m-signing-bonuses/")}),

    ("Opportunity Costs",
     "Return-to-Office: A Pay Cut That Never Hits Payroll",
     [
         "Working from home saves 72 minutes per day on average, across "
         "27 countries (AEA P&P 2023)",
         "Workers put ~40% of the saved time back into the job",
         "Amazon: five days in office from January 2, 2025 (company "
         "announcement)",
         "S&P 500 firms: ~14% higher turnover after RTO mandates, "
         "concentrated among senior staff (working paper)",
     ],
     "Commuting time is compensation priced at zero on the books: a "
     "mandate raises the true cost of the job without touching payroll, "
     "and the market clears it through quits and wage demands. Core "
     "numbers are peer-reviewed (Aksoy, Barrero, Bloom et al.).",
     "If 5 days in office means ~6 unpaid commuting hours a week, what "
     "raise leaves your staff indifferent — and is the in-person gain "
     "bigger than that?",
     "Aksoy et al. (AEA P&P 2023); Barrero, Bloom & Davis (JEP 2023); "
     "aboutamazon.com (Sep 2024); SSRN 5031481",
     {"visual": "72-minute clock split (job / caregiving / leisure) + "
                "commute schematic",
      "notes": (
        "Peer-reviewed anchors: Aksoy, Barrero, Bloom, Davis, Dolls & "
        "Zarate, AEA Papers & Proceedings 113 (2023): WFH saves 72 "
        "min/day, ~40% reinvested in work, ~11% caregiving. Barrero-"
        "Bloom-Davis JEP 2023: full WFH days = 28% of paid US workdays "
        "mid-2023. The 14%-turnover result (Ding et al., SSRN 5031481) "
        "is an unpublished working paper — present as suggestive. Avoid "
        "vendor-survey dollar figures ('$8,158 invisible pay cut').\n"
        "Sources:\n"
        "https://www.aeaweb.org/articles?id=10.1257/pandp.20231013\n"
        "https://www.aeaweb.org/articles?id=10.1257/jep.37.4.23\n"
        "https://www.aboutamazon.com/news/company-news/ceo-andy-jassy-"
        "latest-update-on-amazon-return-to-office-manager-team-ratio\n"
        "https://papers.ssrn.com/sol3/papers.cfm?abstract_id=5031481")}),

    # ---- Sunk costs ---------------------------------------------------------
    ("Sunk Costs",
     "Walking Away from $10 Billion: Apple Car and GM Cruise",
     [
         "Apple (Feb 2024): killed Project Titan after ~10 years, "
         "~2,000 staff, reported >$10B — zero cars shipped",
         "GM (Dec 2024): stopped funding Cruise robotaxis after "
         "reported $10B+ since 2016",
         "GM's forward math (SEC-filed): pay a $0.5B exit charge to "
         "stop a $1.7B-per-year outflow",
         "In both decisions, the $10B already spent appears nowhere",
     ],
     "The modern Concordes, decided correctly: forward margins and "
     "forward costs drove both exits. GM's numbers are filing-confirmed "
     "— a rare fully quantified exit calculation.",
     "GM paid $522M to quit. Under what conditions is paying half a "
     "billion dollars to stop the value-maximizing move?",
     "GM Q4-2024 8-K (SEC); CNBC (Dec 2024 / Feb 2025); NYT & "
     "Bloomberg (Feb 2024; Apple $10B reported, not disclosed)",
     {"visual": "Cruise robotaxi photo + a two-column sunk-vs-forward "
                "ledger for the GM decision",
      "notes": (
        "Confirmed from GM's Q4 2024 8-K: $0.5B Q4 charges (Cruise net "
        "$522M incl. $173M non-cash), $1.7B spent on Cruise in 2024, "
        "~$1B expected annual run-rate savings. Reported (not filed): "
        "the cumulative $10B+ figures for both Apple and GM; Apple has "
        "never disclosed Project Titan spend — say 'reported at more "
        "than $10 billion'. Apple staff largely redirected to the "
        "generative-AI organization; ~1,000 Cruise layoffs Feb 2025, "
        "tech folded into Super Cruise. Honest complication: the Oct "
        "2023 Cruise pedestrian incident changed the FORWARD probability "
        "of scaling — new information, not sunk cost, moved GM.\n"
        "Sources:\n"
        "https://www.sec.gov/Archives/edgar/data/1467858/"
        "000146785825000030/gmq42024pressreleaseandfin.htm\n"
        "https://www.cnbc.com/2024/12/10/gm-halts-funding-of-robotaxi-"
        "development-by-cruise.html\n"
        "https://www.macrumors.com/2024/02/28/apple-car-10-billion-"
        "spent/")}),

    ("Sunk Costs",
     "Meta's Metaverse: $80B of Losses — Patience or Escalation?",
     [
         "Reality Labs operating losses: roughly $80–84B cumulative "
         "2020–2025 (FY2025: $19.2B loss on $2.2B revenue — filings)",
         "Losses rose EVERY year for six years; revenue stuck near $2B",
         "2026 pivot: ~1,500 Reality Labs layoffs; Horizon Worlds moved "
         "to mobile; budget shifted to AI glasses",
         "Meta's defense was forward-looking: option value on the next "
         "computing platform",
     ],
     "The contested case: an escalation-of-commitment pattern versus "
     "genuine platform option value. When Meta finally cut, its stated "
     "reasons were about future returns — not recovering the $80B. "
     "Ideal debate slide after Concorde.",
     "What evidence, visible in 2023 rather than in hindsight, would "
     "distinguish patient platform investment from escalating "
     "commitment?",
     "Meta Q4-2025 results (SEC Ex. 99.1); CNBC (Jan 2026); TechCrunch "
     "(Feb 2026)",
     {"visual": "Native bar chart: Reality Labs annual operating losses "
                "2020–2025",
      "notes": (
        "Confirmed from filings/tier-1: FY2025 loss $19.19B on $2.21B "
        "revenue; FY2024 $17.73B; Q4-2025 $6.02B. Cumulative 2020–25 "
        "sums to ~$83.6B from disclosed segment losses (2020 $6.62B, "
        "2021 $10.19B, 2022 $13.71B, 2023 $16.12B) — say 'roughly "
        "$80–84B'. Reported: ~1,500 layoffs (~10%) Jan 2026; Horizon "
        "Worlds to mobile 20 Feb 2026; 'up to 30%' metaverse budget cut "
        "is press-reported only.\n"
        "Sources:\n"
        "https://www.sec.gov/Archives/edgar/data/1326801/"
        "000162828026003832/meta-12312025xexhibit991.htm\n"
        "https://www.cnbc.com/2026/01/28/metas-reality-labs-posts-"
        "6point02-billion-loss-in-fourth-quarter.html\n"
        "https://techcrunch.com/2026/02/20/meta-metaverse-leaves-vr-"
        "horizon-worlds-mobile/")}),

    # ---- Cost-benefit and marginal analysis --------------------------------
    ("Cost-Benefit and Marginal Analysis",
     "Cutting the Marginal Flight: Airlines and the 2026 Fuel Shock",
     [
         "Jet fuel “more than doubled in three weeks” (Mar 2026); "
         "United cut ~5% of planned capacity",
         "Cuts hit off-peak and red-eyes — flights “15% to 20% less "
         "valuable” than average (Delta's CCO)",
         "Kirby: “no point in burning cash in the near term… nothing "
         "changes about our longer-term plans”",
         "Fleet and 2027 capacity untouched: only the flights below the "
         "new MB=MC crossing were cut",
     ],
     "MB = MC with the marginal unit named by the CEO: when marginal "
     "cost jumps, you shed exactly the units below the new crossing — "
     "the marginal flight, not the average one. Lands right after the "
     "exercise-hours example.",
     "Which costs entered the flight-cancellation decision and which "
     "did not — and would you have made the same split?",
     "Kirby staff memo via Fox Business / CNBC (Mar 2026); Travel "
     "Market Report (Apr 2026, Delta)",
     {"visual": "Native MB/MC chart with the crossing shifting left + "
                "red-eye flight schedule graphic",
      "notes": (
        "Verified from coverage of the 20 Mar 2026 Kirby memo: ~5% "
        "capacity cut (≈3pp off-peak/midweek/overnight, ~1pp O'Hare, "
        "~1pp Tel Aviv/Dubai suspensions); modeling oil to $175/bbl and "
        ">$100 through 2027 (≈$11B added fuel bill). Delta (Apr 2026, "
        "EVP Joe Esposito): off-peak/edge-of-day/red-eyes '15–20% less "
        "valuable on a net revenue basis'. CNBC returned 403 — its "
        "detail is via search snippet; Fox Business and Travel Market "
        "Report fetched directly.\n"
        "Sources:\n"
        "https://www.foxbusiness.com/economy/united-airlines-slashes-"
        "flights-as-iran-war-sends-fuel-prices-soaring\n"
        "https://www.travelmarketreport.com/air/articles/delta-air-"
        "lines-to-cut-capacity-due-to-rising-oil-prices")}),

    ("Market Mechanism and Fairness",
     "After the LA Fires: A 10% Rent Cap Meets a Housing Shortage",
     [
         "Jan 7, 2025 emergency: CA Penal Code §396 capped rent "
         "increases at 10% — with criminal penalties",
         "Illegal to ACCEPT more — even if the tenant volunteers it "
         "(AG Bonta)",
         "750+ warning letters by Jun 2025; realtors criminally charged "
         "(one listing >50% above pre-fire rent)",
         "Scraped Zillow data (advocacy group, reported): ~15% of "
         "updated listings above the cap; 2-bedrooms up ~80%",
     ],
     "The “shortages when disasters loom” discussion (slide 29) with a "
     "case from the students' own zip codes: a binding cap converts a "
     "price rise into non-price rationing — queues, networks, "
     "non-compliance — and writes the fairness intuition into criminal "
     "law.",
     "If the cap had been lifted on Jan 8 and rents had doubled — who "
     "gets the housing, and is that allocation better or worse?",
     "CA AG press releases (Jan–Jun 2025); Gov. Newsom EOs (Mar 2025, "
     "Jan 2026); NLIHC / Rent Brigade (reported)",
     {"visual": "LA-fires photo + native S/D chart with a price ceiling "
                "below the market-clearing rent",
      "notes": (
        "Confirmed from CA AG / Governor's office: §396 activated 7 Jan "
        "2025; 10% cap; penalties up to 1 year jail / $10,000 + $2,500 "
        "civil per violation; 200+ warning letters by 18 Jan, 650+ by "
        "28 Jan (second realtor charged, Glendale listing >50% above "
        "list), 750+ by Jun 2025; protections extended to 1 Jul 2025. "
        "Reported (advocacy scrape, not official): Rent Brigade Zillow "
        "analysis — 15% of updated listings (1,343) above cap, studios "
        "+57%, 2BRs +80%, only 9% corrected. OPEN: whether a housing-"
        "specific cap is still active in LA County as of Aug 2026 (the "
        "Jan 2026 EO covers building materials/reconstruction).\n"
        "Sources:\n"
        "https://www.oag.ca.gov/news/press-releases/attorney-general-"
        "bonta-rental-bidding-wars-during-la-wildfires-prohibited-under\n"
        "https://www.oag.ca.gov/news/press-releases/attorney-general-"
        "bonta-charges-second-los-angeles-realtor-price-gouging-victims\n"
        "https://nlihc.org/resource/report-finds-emergency-declaration-"
        "los-angeles-wildfires-was-followed-rent-gouging")}),
]


# --------------------------------------------------------------------------
# Bench slides — runners-up, one line each (promotable on request)
# --------------------------------------------------------------------------
BENCH_1 = [
    ("Market definition", [
        "FTC v. Meta (Nov 2025): narrow market REJECTED — TikTok/YouTube "
        "converged; the counterexample to Tapestry/Kroger",
        "Google search remedies (Sep 2025): are AI chatbots in the "
        "market? Outside for liability, inside for the remedy",
        "FTC v. Amazon “online superstore”: two adjectives do all the "
        "work; trial Mar 2027 — prompt, not case",
    ]),
    ("Supply and demand", [
        "US beef 2024–26: +37% (BLS), smallest herd since 1951 — slow "
        "biological supply response vs. eggs' fast one",
        "Coffee: retail +52% (BLS-verified) while futures retreated — "
        "the pass-through lag executives misread",
        "Cocoa: price fell ~70% while consumption ALSO fell — you need "
        "quantity, not just price, to identify the shift",
        "GLP-1 drugs: administered prices → queues, waitlists, and a "
        "compounding gray market instead of price increases",
    ]),
]

BENCH_2 = [
    ("Opportunity costs", [
        "Berkshire's record ~$397B cash (Q1 2026, reported) and Abel's "
        "deployments: OxyChem $9.5B, Taylor Morrison, buybacks",
        "Hyperscaler AI capex ~$700B (2026, big four): the required "
        "return that never appears on the income statement",
        "NIL-era draft decisions: staying in college now has a price "
        "tag ($6.5M reported NIL offer vs. NFL rookie deal)",
    ]),
    ("Sunk costs", [
        "California high-speed rail: approved 2008 at $9B, Phase I now "
        "$126B (LAO) — the fallacy in its political habitat",
        "Ørsted's Hornsea 4 (May 2025): paid up to ~DKK 5B to quit "
        "despite a locked-in price — quitting has a price tag",
        "Sony's Concord (2024): finished game taken offline in 14 days "
        "with full refunds — “it's finished” is not a reason",
    ]),
]

BENCH_3 = [
    ("Marginal analysis and fairness", [
        "Google datacenter demand response (1 GW by 2026): sell the "
        "megawatt, not the compute — the marginal-value comparison",
        "Wendy's “surge pricing” (2024): identical mechanism, opposite "
        "reception as discount vs. surcharge — the framing lesson",
        "Egg rationing 2025: Trader Joe's quantity limits vs. Waffle "
        "House's surcharge; DOJ 2026 benchmark-manipulation settlement",
    ]),
]


def bench_slide(prs, page_num, title, groups):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, "Module 1 · Candidates · Bench")
    _draw_action_title(slide, title)
    items = []
    for concept, lines in groups:
        items.append((concept, 0, {'bold': True, 'bullet_style': 'none',
                                   'size': 20}))
        for ln in lines:
            items.append((ln, 1, {'size': 18}))
    box = _add_hierarchical_bullets(
        slide, left=MARGIN, top=Inches(1.55), width=RULE_W,
        height=Inches(5.4), items=items,
        size=20, sub_size=18, line_spacing_pts=8, sub_line_spacing_pts=3)
    box.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    _footer(slide, page_num)
    return slide


def build(out_path=None):
    prs = Presentation()
    prs.slide_width = int(SLIDE_W)
    prs.slide_height = int(SLIDE_H)
    slide_cover(prs)                                       #  1
    # Expanded (adoption-ready) examples replace CANDIDATES[0..2] and
    # CANDIDATES[9], which stay in the list above only as the original
    # one-slide summaries (not built).
    exp_tapestry_case(prs, 2)                              #  2
    exp_tapestry_evidence(prs, 3)                          #  3
    exp_kroger_case(prs, 4)                                #  4
    exp_costco_run(prs, 5)                                 #  5
    exp_netflix_saga(prs, 6)                               #  6
    exp_netflix_chart(prs, 7)                              #  7
    c = CANDIDATES[3]              # eggs
    candidate_slide(prs, 8, *c[:6], **(c[6] if len(c) > 6 else {}))
    # DRAM expanded to 3 slides (replaces CANDIDATES[4]):
    exp_dram_case(prs, 9)                                  #  9
    exp_dram_sd(prs, 10)                                   # 10
    exp_dram_outcome(prs, 11)                              # 11
    n = 12
    for c in CANDIDATES[5:9]:      # talent, RTO, AppleGM, Meta
        extras = c[6] if len(c) > 6 else {}
        candidate_slide(prs, n, *c[:6], **extras)
        n += 1
    exp_united_case(prs, n)                                # 16
    exp_united_chart(prs, n + 1)                           # 17
    c = CANDIDATES[10]             # LA rent cap
    candidate_slide(prs, n + 2, *c[:6], **(c[6] if len(c) > 6 else {}))
    bench_slide(prs, n + 3, "Bench: More Candidates I — Markets, Supply "
                            "and Demand", BENCH_1)
    bench_slide(prs, n + 4, "Bench: More Candidates II — Opportunity "
                            "and Sunk Costs", BENCH_2)
    bench_slide(prs, n + 5, "Bench: More Candidates III — Marginal "
                            "Analysis and Fairness", BENCH_3)
    out = Path(out_path) if out_path else \
        OUT_DIR / "Module 1 - Example Candidates.pptx"
    prs.save(str(out))
    print(f"saved {out} — {len(prs.slides._sldIdLst)} slides")
    return out


if __name__ == "__main__":
    import sys as _sys
    build(_sys.argv[1] if len(_sys.argv) > 1 else None)
