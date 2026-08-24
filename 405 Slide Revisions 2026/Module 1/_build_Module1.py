# ==========================================================================
#  _build_Module2InClass.py — phase-1 scaffold for
#  "Module 2 - In Class Revised.pptx"
#
#  Demand Analysis, in-class deck (76 slides per the approved outline in
#  "Module 2 - In Class Revised - outline.md"; the 16 PollEverywhere
#  slides + the pizza Excel-embed slide are positional stubs until the
#  _splice_media.py pass).
#
#  Helper layer copied VERBATIM from Module 7/_build_Module7.py
#  (2026-08-14), which itself carries the Module 3 helper layer — the
#  proven chrome, box, bullet, OMML, chart, and table primitives.
#  M2-specific code starts at the "MODULE 2" banner below.
# ==========================================================================

import copy
import math
import re
import shutil
import zipfile
from pathlib import Path

from lxml import etree as ET
from pptx import Presentation
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn
from pptx.util import Cm, Inches, Pt

# Reuse all primitives from the template script (single source of truth).
from _build_template_samples import (
    FADED,
    FOOTER_TEXT,
    GOLD,
    GOLD_W,
    GRAY,
    MARGIN,
    MODULE_AGENDA,
    NAVY,
    RULE,
    RULE_W,
    SLIDE_H,
    SLIDE_W,
    WHITE,
    _add_bulleted_list,
    _add_rect,
    _add_text,
    _blank_slide,
    _draw_action_title,
    _set_bullet_char,
)

OUT_DIR = Path(__file__).parent


# --------------------------------------------------------------------------
# Title-case top bar – replaces the all-caps default from the template script.
# The user prefers academic-paper title case (each major word capitalized,
# small connectors like "and", "of", "for", "the" stay lowercase).
# --------------------------------------------------------------------------

def _draw_top_bar_tc(slide, section_tag):
    """Navy top bar with title-cased section tag (no uppercase forcing)."""
    bar_h = Inches(0.42)
    _add_rect(slide, 0, 0, SLIDE_W, bar_h, NAVY)
    _add_text(slide, MARGIN, 0, Inches(12), bar_h,
              section_tag, size=16, bold=True,
              color=WHITE, font="Calibri",
              anchor=MSO_ANCHOR.MIDDLE)


def _draw_footer(slide, footer_text, page_num):
    """Footer rule + gold accent + footer text/page number, with larger type
    than the template default so handout printouts remain legible."""
    _add_rect(slide, 0, Inches(7.15), SLIDE_W, Inches(0.02), RULE)
    _add_rect(slide, MARGIN, Inches(7.135), GOLD_W, Inches(0.05), GOLD)
    _add_text(slide, MARGIN, Inches(7.20), Inches(11), Inches(0.32),
              footer_text, size=12, color=GRAY)
    _add_text(slide, Inches(12.5), Inches(7.20), Inches(0.6), Inches(0.32),
              str(page_num), size=12, color=GRAY, align=PP_ALIGN.RIGHT)


# --------------------------------------------------------------------------
# Speaker-notes helper
# --------------------------------------------------------------------------

def _set_notes(slide, text):
    """Replace the slide's speaker notes with *text*."""
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.clear()
    notes_tf.text = text


# --------------------------------------------------------------------------
# Reusable shape primitives for diagrams
# --------------------------------------------------------------------------

def _add_filled_box(slide, left, top, width, height, label, *,
                    fill=NAVY, text_color=WHITE, line=None,
                    size=18, bold=True, font="Calibri"):
    """Filled rectangle with centered text."""
    left, top, width, height = int(left), int(top), int(width), int(height)
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height,
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = text_color
    return shp


def _add_outlined_box(slide, left, top, width, height, label, *,
                      line=NAVY, text_color=NAVY, fill=WHITE,
                      size=18, bold=True, line_w=1.25, font="Calibri",
                      rounded=False, shadow=False, corner_pct=0.06):
    """Outlined rectangle (white fill) with centered text.

    ``rounded=True`` switches the base shape to a rounded rectangle
    (corner-adjust = ``corner_pct``); ``shadow=True`` adds a soft drop
    shadow.  Both default off so existing flat call-sites are
    unaffected.
    """
    left, top, width, height = int(left), int(top), int(width), int(height)
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(
        shape_type, left, top, width, height,
    )
    if rounded:
        try: shp.adjustments[0] = corner_pct
        except Exception: pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    if shadow:
        _add_drop_shadow(shp)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = text_color
    return shp


def _add_convention_box(slide, left, top, width, height, *,
                          prefix=None, body=None, runs=None,
                          fill_rgb=None, border=None, line_w=1.0,
                          corner_pct=0.12, size=15, align=PP_ALIGN.LEFT,
                          font="Calibri", pad_h=None, pad_v=None,
                          line_spacing_pct=None):
    """Cream-fill / navy-border rounded-rect explanation callout.

    The "Convention" textbox pattern from slide 14 generalised — use it
    anywhere a slide needs a compact, visually-distinct box for a
    short conceptual explanation or notational convention.  Sits well
    below a table, beside a hero formula, or as a slide-wide footer.

    Two ways to populate the text:
      • ``prefix`` (bold) + ``body`` (regular) — simplest path; matches
        slide 14's "Convention:  <text>" pattern.
      • ``runs`` — a list of ``(text, {"bold": .., "italic": .., ...})``
        tuples for finer-grained styling (multi-line, mixed formatting).

    Style defaults follow the course-layer CLAUDE.md "Convention callout
    box" spec — cream fill, thin primary-color border, slight rounding,
    primary-color text.  Override ``fill_rgb`` / ``border`` only when you
    need a different accent.
    """
    fill = fill_rgb if fill_rgb is not None else RGBColor(0xFD, 0xF6, 0xE6)
    border = border if border is not None else NAVY

    left, top, width, height = int(left), int(top), int(width), int(height)
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border
    shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    try:
        shp.adjustments[0] = corner_pct
    except Exception:
        pass

    # Inset text box so the rounded corners breathe — matches slide 14.
    pad_h = Inches(0.20) if pad_h is None else pad_h
    pad_v = Inches(0.12) if pad_v is None else pad_v
    tb = slide.shapes.add_textbox(
        left + int(pad_h), top + int(pad_v),
        width - 2 * int(pad_h), height - 2 * int(pad_v),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0); tf.margin_bottom = Inches(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    def _style_run(r, opts):
        r.font.name = opts.get('font', font)
        r.font.size = Pt(opts.get('size', size))
        r.font.bold = opts.get('bold', False)
        r.font.italic = opts.get('italic', False)
        r.font.underline = opts.get('underline', False)
        r.font.color.rgb = opts.get('color', NAVY)

    if runs is not None:
        first = True
        for entry in runs:
            text, opts = entry if isinstance(entry, tuple) else (entry, {})
            if opts.get('newline') and not first:
                p = tf.add_paragraph()
            elif first:
                p = tf.paragraphs[0]
            else:
                # Same paragraph — append run to the most-recent paragraph.
                p = tf.paragraphs[-1]
            p.alignment = align
            r = p.add_run()
            r.text = text
            _style_run(r, opts)
            first = False
    else:
        p = tf.paragraphs[0]
        p.alignment = align
        if prefix:
            r1 = p.add_run(); r1.text = prefix
            _style_run(r1, {'bold': True, 'color': NAVY, 'size': size})
        if body:
            r2 = p.add_run(); r2.text = body
            _style_run(r2, {'color': NAVY, 'size': size})

    if line_spacing_pct is not None:
        for p_obj in tf.paragraphs:
            pPr = p_obj._p.get_or_add_pPr()
            for old in pPr.findall(qn('a:lnSpc')):
                pPr.remove(old)
            lnSpc = ET.Element(qn('a:lnSpc'))
            spcPct = ET.SubElement(lnSpc, qn('a:spcPct'))
            spcPct.set('val', str(int(line_spacing_pct * 1000)))
            pPr.insert(0, lnSpc)
    return shp


def _add_rounded_filled_box(slide, left, top, width, height, label, *,
                             fill=NAVY, text_color=WHITE, line=None,
                             size=18, bold=True, font="Calibri",
                             corner_pct=0.06, shadow=True):
    """Rounded-corner filled rectangle with centered text and soft drop shadow.

    Mirrors :func:`_add_filled_box` but renders ``MSO_SHAPE.ROUNDED_RECTANGLE``
    with the corner-adjust set to ``corner_pct`` (6 % per course CLAUDE.md
    "slight rounding") and a soft drop shadow via :func:`_add_drop_shadow`.
    """
    left, top, width, height = int(left), int(top), int(width), int(height)
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
    )
    try:
        shp.adjustments[0] = corner_pct
    except Exception:
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    if shadow:
        _add_drop_shadow(shp)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = text_color
    return shp


def _add_arrow(slide, start_xy, end_xy, *, color=NAVY, weight_pt=1.5,
               head=True, dash=None, head_size='med'):
    """Draw a line/arrow from start to end (in EMU/Inches values).

    EMU coordinates MUST be integers — PowerPoint rejects decimal values
    in <a:off>/<a:ext> and refuses to open the file. Cast to int defensively.

    ``dash`` accepts any OOXML preset-dash name (e.g., ``"dash"``,
    ``"dashDot"``, ``"sysDash"``).  Default ``None`` = solid line.

    ``head_size`` is the OOXML preset arrowhead size — one of ``'sm'``,
    ``'med'`` (default), or ``'lg'``.  Width and height are set
    together, so passing ``'lg'`` gives a noticeably larger tip while
    leaving the line weight unchanged.
    """
    sx, sy = int(start_xy[0]), int(start_xy[1])
    ex, ey = int(end_xy[0]), int(end_xy[1])
    line = slide.shapes.add_connector(1, sx, sy, ex, ey)  # 1 = STRAIGHT
    line.line.color.rgb = color
    line.line.width = Pt(weight_pt)
    ln = line.line._get_or_add_ln()
    if dash is not None:
        for old in ln.findall(qn('a:prstDash')):
            ln.remove(old)
        prst = ET.SubElement(ln, qn('a:prstDash'))
        prst.set('val', dash)
    if head:
        tailEnd = ET.SubElement(ln, qn('a:tailEnd'))
        tailEnd.set('type', 'triangle')
        tailEnd.set('w', head_size)
        tailEnd.set('h', head_size)
    return line


def _add_wavy_line(slide, x_start, x_end, y_center, *,
                    amplitude=None, cycles=1.75, segments=36,
                    color=NAVY, weight_pt=1.5):
    """Horizontal sinusoidal line from x_start to x_end at y_center.

    Renders a polyline approximation of ``sin(2π · t · cycles)`` inside a
    custGeom shape so the line reads as a gentle wave rather than a
    straight connector.  ``amplitude`` is the peak-to-baseline height in
    EMU; defaults to ~0.04".  ``segments`` controls how finely the wave
    is discretised — 30+ is smooth enough that the polyline reads as a
    curve.  No arrowhead.
    """
    if amplitude is None:
        amplitude = Inches(0.04)
    L = int(x_end - x_start)
    A = int(amplitude)
    if L == 0:
        return None
    flip = "1" if L < 0 else "0"
    bbox_left = int(min(x_start, x_end))
    bbox_top = int(y_center - A)
    bbox_w = abs(L)
    bbox_h = 2 * A

    pts = []
    for i in range(segments + 1):
        t = i / segments
        lx = int(round(t * 100000))
        sin_val = math.sin(2 * math.pi * t * cycles)
        # Path coords: y=0 is top.  Centre at 50000; +1·amplitude → top
        # (0), −1·amplitude → bottom (100000).
        ly = int(round(50000 - sin_val * 50000))
        pts.append((lx, ly))

    path_segs = [f'<a:moveTo><a:pt x="{pts[0][0]}" y="{pts[0][1]}"/></a:moveTo>']
    for lx, ly in pts[1:]:
        path_segs.append(f'<a:lnTo><a:pt x="{lx}" y="{ly}"/></a:lnTo>')
    path_inner = ''.join(path_segs)

    color_hex = f'{color[0]:02X}{color[1]:02X}{color[2]:02X}'
    width_emu = int(weight_pt * 12700)

    P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    A_NS_LOCAL = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    sp_xml = (
        f'<p:sp xmlns:p="{P_NS}" xmlns:a="{A_NS_LOCAL}">'
        f'<p:nvSpPr><p:cNvPr id="0" name="WavyLine"/>'
        f'<p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr>'
        f'<a:xfrm flipH="{flip}">'
        f'<a:off x="{bbox_left}" y="{bbox_top}"/>'
        f'<a:ext cx="{bbox_w}" cy="{bbox_h}"/>'
        f'</a:xfrm>'
        f'<a:custGeom>'
        f'<a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/>'
        f'<a:rect l="0" t="0" r="0" b="0"/>'
        f'<a:pathLst>'
        f'<a:path w="100000" h="100000" fill="none">'
        f'{path_inner}'
        f'</a:path>'
        f'</a:pathLst>'
        f'</a:custGeom>'
        f'<a:noFill/>'
        f'<a:ln w="{width_emu}" cap="rnd">'
        f'<a:solidFill><a:srgbClr val="{color_hex}"/></a:solidFill>'
        f'</a:ln>'
        f'</p:spPr>'
        f'</p:sp>'
    )
    elem = ET.fromstring(sp_xml)
    slide.shapes._spTree.append(elem)
    return elem


def _add_arrow_shape(slide, left, top, width, height, *,
                     direction="right", fill=GOLD, line=None):
    """Block arrow shape (the 'we-are-here' indicator).

    direction: "right" (default), "left", "up", or "down".
    """
    left, top, width, height = int(left), int(top), int(width), int(height)
    geom_map = {
        "left": MSO_SHAPE.LEFT_ARROW,
        "right": MSO_SHAPE.RIGHT_ARROW,
        "up": MSO_SHAPE.UP_ARROW,
        "down": MSO_SHAPE.DOWN_ARROW,
    }
    geom = geom_map.get(direction, MSO_SHAPE.RIGHT_ARROW)
    shp = slide.shapes.add_shape(geom, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp



# --------------------------------------------------------------------------
# Layout 2 — Section Header / Agenda (parameterized)
# --------------------------------------------------------------------------

def make_section_agenda(prs, page_num, *, current_part_idx=None,
                        current_sub_idx=None,
                        section_tag="Module 3 · Agenda",
                        title="Agenda"):
    """Render an agenda / section-divider slide.

    ``current_part_idx`` makes that Part navy and fades the others.
    ``current_sub_idx`` (optional, only meaningful with current_part_idx)
    further restricts the navy highlight inside the current Part to a
    single sub-bullet — the other subs render in faded gray, alongside
    Part 2.  Used by intra-Part dividers (e.g., the new Short Run agenda
    on page 13, or the Long Run agenda on page 31).
    """
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, section_tag)
    _draw_action_title(slide, title)

    y = Inches(1.85)
    part_title_h = Inches(0.6)
    sub_list_h = Inches(1.35)
    block_gap = Inches(0.15)

    for idx, part in enumerate(MODULE_AGENDA):
        if current_part_idx is None:
            color = NAVY  # preview: highlight all Parts
        else:
            color = NAVY if idx == current_part_idx else FADED

        # 2026-05-19: per-sub fade.  When this is the current Part AND a
        # specific sub-index is highlighted, give every other sub the
        # FADED color so the deck shows "you are HERE" within the Part.
        if (current_part_idx is not None and idx == current_part_idx
                and current_sub_idx is not None):
            sub_colors = [
                NAVY if i == current_sub_idx else FADED
                for i in range(len(part["subs"]))
            ]
        else:
            sub_colors = None

        _add_text(slide, MARGIN, y, RULE_W, part_title_h,
                  part["title"], size=30, bold=True, color=color,
                  font="Calibri")
        y += part_title_h

        _add_bulleted_list(
            slide,
            left=MARGIN + Inches(0.4),
            top=y,
            width=RULE_W - Inches(0.4),
            height=sub_list_h,
            items=part["subs"],
            size=24, color=color, bullet_color=color,
            # 2026-05-19: sub-bullets tightened from 10 → 3 pt to match
            # the deck-wide rhythm (sub-bullets cluster under their
            # parent rather than breathing out).  Applied to both
            # Parts on every agenda divider in the deck.
            line_spacing_pts=3,
            autonum_scheme='alphaLcPeriod',
            colors=sub_colors,
        )
        y += sub_list_h + block_gap

    _draw_footer(slide, FOOTER_TEXT, page_num)
    return slide


# --------------------------------------------------------------------------
# Layout 3 — Content bulleted
# --------------------------------------------------------------------------

def make_content_bulleted(prs, page_num, section_tag, title, bullets, *,
                          size=24, sub_size=None, line_spacing_pts=18,
                          sub_line_spacing_pts=None,
                          extras=None, bullets_top=None):
    """bullets: list of (text, level) tuples OR plain strings (level=0).

    ``bullets_top`` overrides the default body-region start (Inches(1.85))
    — useful when the slide also hosts large diagrams below the bullets
    and the bullets need to lift up to avoid overlap.

    ``sub_line_spacing_pts`` overrides the legacy
    ``max(6, line_spacing_pts - 8)`` formula for sub-bullet space-before.
    """
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, section_tag)
    _draw_action_title(slide, title)

    normalized = [(b, 0) if isinstance(b, str) else b for b in bullets]

    if bullets_top is None:
        bullets_top = Inches(1.85)

    _add_hierarchical_bullets(
        slide,
        left=MARGIN,
        top=bullets_top,
        width=RULE_W,
        height=Inches(5.0),
        items=normalized,
        size=size,
        sub_size=sub_size,
        line_spacing_pts=line_spacing_pts,
        sub_line_spacing_pts=sub_line_spacing_pts,
    )

    if extras is not None:
        extras(slide)

    _draw_footer(slide, FOOTER_TEXT, page_num)
    return slide


def _inject_bullet_lst_style(tf, *, size, sub_size,
                              main_color=NAVY, sub_color=GRAY,
                              main_char='▪', sub_char='–',
                              main_space_pts=None, sub_space_pts=None):
    """Inject <a:lstStyle> into a text frame defining per-level defaults.

    Defines lvl1pPr (main bullets, lvl="0") and lvl2pPr (sub bullets,
    lvl="1") with bullet character, indent, color, font, and space-
    before defaults.  Enables PowerPoint's native Tab / Shift+Tab to
    auto-reformat a bullet when its level changes: runs / paragraphs
    that do NOT explicitly override these properties inherit them.

    main_space_pts / sub_space_pts: if set, the <a:spcBef> for each
    level.  Sub-bullets default to a tighter spacing than main bullets
    (3 pt vs ~10 pt) per the deck's pedagogical preference: sub-bullets
    visually cluster under their parent, main bullets give the eye a
    larger break.
    """
    txBody = tf._txBody
    # Remove any prior lstStyle
    for old in txBody.findall(qn('a:lstStyle')):
        txBody.remove(old)
    lstStyle = ET.Element(qn('a:lstStyle'))
    levels = [
        ('a:lvl1pPr', 342900, -342900, main_color, main_char, size, main_space_pts),
        ('a:lvl2pPr', 685800, -228600, sub_color, sub_char, sub_size, sub_space_pts),
    ]
    for tag, mar_l, indent, color, char, sz, space_pts in levels:
        lvl = ET.SubElement(lstStyle, qn(tag))
        lvl.set('marL', str(mar_l))
        lvl.set('indent', str(indent))
        # spcBef must precede bullet attributes (OOXML schema order).
        if space_pts is not None:
            spcBef = ET.SubElement(lvl, qn('a:spcBef'))
            spcPts = ET.SubElement(spcBef, qn('a:spcPts'))
            spcPts.set('val', str(int(space_pts * 100)))
        bc = ET.SubElement(lvl, qn('a:buClr'))
        sc = ET.SubElement(bc, qn('a:srgbClr'))
        sc.set('val', '{:02X}{:02X}{:02X}'.format(color[0], color[1], color[2]))
        bf = ET.SubElement(lvl, qn('a:buFont'))
        bf.set('typeface', 'Calibri')
        bch = ET.SubElement(lvl, qn('a:buChar'))
        bch.set('char', char)
        drp = ET.SubElement(lvl, qn('a:defRPr'))
        drp.set('sz', str(int(sz * 100)))
        drp.set('b', '0')
        sf = ET.SubElement(drp, qn('a:solidFill'))
        sfc = ET.SubElement(sf, qn('a:srgbClr'))
        sfc.set('val', '{:02X}{:02X}{:02X}'.format(color[0], color[1], color[2]))
        lt = ET.SubElement(drp, qn('a:latin'))
        lt.set('typeface', 'Calibri')
    # Insert lstStyle after bodyPr (proper element order in <a:txBody>)
    bodyPr = txBody.find(qn('a:bodyPr'))
    if bodyPr is not None:
        bodyPr.addnext(lstStyle)
    else:
        txBody.insert(0, lstStyle)


def _add_hierarchical_bullets(slide, left, top, width, height, items,
                              *, size=24, sub_size=None, line_spacing_pts=18,
                              sub_line_spacing_pts=None):
    """Render bullets with indent levels.

    Bullet item forms:
        (text, level)                       — simple, defaults from level
        (text, level, opts)                 — opts dict overrides
        (runs_list, level, opts)            — multi-run paragraph

    text:
        - str  → a single run with text
        - ''   → empty paragraph (visual spacer; no run)
        - list → multi-run: list of ``(run_text, run_opts)`` tuples

    Paragraph-level opts (all optional):
        bullet_style: 'main' (▪ NAVY) | 'sub' (– GRAY) | 'arrow' (no bullet,
            plain left-indent; the run text supplies the leader char like
            "→" or Wingdings) | 'none'
        mar_l, indent: bullet positioning (EMU)
        space_before_pts: spcBef in pts (overrides legacy formula)
        size, color, bold, italic: defaults applied to every run unless
            the run_opts override them.

    Run-level opts (run_opts in a runs_list tuple):
        font_name (default 'Calibri'), size, color, bold, italic,
        underline (bool), wingdings (bool — emits <a:sym typeface="Wingdings"/>
        so a private-use-area character renders as its Wingdings glyph).

    Each paragraph also receives a ``lvl="N"`` attribute when level > 0
    so PowerPoint's Tab / Shift-Tab outline navigation can find an
    explicit outline level.
    """
    if sub_size is None:
        sub_size = size - 4

    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    # 2026-05-18: inject <a:lstStyle> with per-level defaults so that
    # PowerPoint's Tab / Shift+Tab keyboard shortcuts auto-reformat the
    # bullet when its outline level changes.  Default 'main'/'sub'
    # bullets (no custom indent, level < 2) skip explicit pPr/run
    # styling — they inherit from lstStyle.  Bullets with custom indent
    # OR level >= 2 OR explicit size/color overrides still apply
    # explicit styling.
    # 2026-05-18 (later): also moved space-before into lstStyle so a
    # paragraph demoted via Tab picks up the sub-bullet's tighter spcBef
    # automatically.  Sub-bullet default: 3 pt (was 6 pt) — per user
    # preference, sub-bullets cluster tightly under their parent.
    sub_default_space = (sub_line_spacing_pts if sub_line_spacing_pts is not None
                          else 3)
    _inject_bullet_lst_style(tf, size=size, sub_size=sub_size,
                              main_space_pts=line_spacing_pts,
                              sub_space_pts=sub_default_space)

    for i, item in enumerate(items):
        if len(item) == 2:
            text, level = item
            opts = {}
        else:
            text, level, opts = item
            opts = opts or {}

        # Determine inheritance up-front so spcBef logic below can use it.
        style = opts.get('bullet_style', 'main' if level == 0 else 'sub')
        has_custom_indent = 'mar_l' in opts or 'indent' in opts
        inherits_from_lst = (style in ('main', 'sub')
                              and not has_custom_indent
                              and level < 2)

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        pPr = p._p.get_or_add_pPr()

        # Space-before.  When the paragraph inherits from lstStyle and the
        # caller has not overridden, skip the per-pPr spcBef so that Tab-
        # induced level changes pick up the new level's spcBef from
        # lstStyle.  Otherwise set explicitly.
        if i > 0:
            sp_override = opts.get('space_before_pts')
            if sp_override is not None:
                spcBef = ET.SubElement(pPr, qn('a:spcBef'))
                pts = ET.SubElement(spcBef, qn('a:spcPts'))
                pts.set('val', str(sp_override * 100))
            elif not inherits_from_lst:
                if level == 0:
                    sp = line_spacing_pts
                else:
                    sp = sub_default_space
                spcBef = ET.SubElement(pPr, qn('a:spcBef'))
                pts = ET.SubElement(spcBef, qn('a:spcPts'))
                pts.set('val', str(sp * 100))
            # else: inherit from lstStyle.

        # Outline-level attribute (enables PowerPoint Tab/Shift-Tab)
        if level > 0:
            pPr.set('lvl', str(level))

        # Bullet styling — default 'main' (lvl 0) and 'sub' (lvl 1)
        # inherit from lstStyle when no custom indent is requested.
        # 'arrow' and 'none' suppress the lstStyle bullet via <a:buNone/>.
        # Level >= 2 still uses explicit _set_bullet_char (lstStyle here
        # only defines lvl1pPr and lvl2pPr).
        if style == 'main':
            if not inherits_from_lst:
                _set_bullet_char(p, char='▪', color=NAVY,
                                  mar_l=opts.get('mar_l', 342900),
                                  indent=opts.get('indent', -342900),
                                  size_pct=100)
        elif style == 'sub':
            if not inherits_from_lst:
                default_mar = 342900 + level * 342900
                _set_bullet_char(p, char='–', color=GRAY,
                                  mar_l=opts.get('mar_l', default_mar),
                                  indent=opts.get('indent', -228600),
                                  size_pct=100)
        elif style == 'arrow':
            # Plain left-indent, no bullet glyph; user text supplies leader.
            pPr.set('marL', str(opts.get('mar_l', 457200)))
            if 'indent' in opts:
                pPr.set('indent', str(opts['indent']))
            # Explicit <a:buNone/> so the lstStyle bullet doesn't bleed
            # through.
            ET.SubElement(pPr, qn('a:buNone'))
        elif style == 'none':
            ET.SubElement(pPr, qn('a:buNone'))

        # Empty paragraph (spacer) — no run; suppress the bullet so the
        # empty line doesn't render an orphan glyph.
        if text == '':
            if inherits_from_lst:
                ET.SubElement(pPr, qn('a:buNone'))
            continue

        # Normalize text to a runs list
        if isinstance(text, str):
            runs = [(text, {})]
        else:
            runs = text  # already a list of (run_text, run_opts) tuples

        # Paragraph-level explicit overrides (None = inherit from lstStyle).
        para_size_override = opts.get('size')
        para_color_override = opts.get('color')
        para_bold_override = opts.get('bold')
        para_italic_override = opts.get('italic')

        for run_text, run_opts in runs:
            run_opts = run_opts or {}

            # Inline OMML math zone — append <a14:m><m:oMath>...</m:oMath></a14:m>
            # as a sibling of the regular <a:r> runs so the formula renders
            # in-line with surrounding prose.  ``run_text`` is interpreted
            # as raw OMML content (e.g., from `_formula_mp_ratio` or the
            # `_omml_*` builders).
            if run_opts.get('omml'):
                om_sz = run_opts.get('size')
                if om_sz is None:
                    om_sz = para_size_override
                if om_sz is None:
                    om_sz = size if level == 0 else sub_size
                om_clr = run_opts.get('color') or para_color_override
                if om_clr is None:
                    om_clr = NAVY if level == 0 else GRAY
                sz_centi = int(om_sz * 100)
                clr_hex = '{:02X}{:02X}{:02X}'.format(
                    om_clr[0], om_clr[1], om_clr[2])
                a14_xml = (
                    f'<a14:m xmlns:a14="{A14_NS}" '
                    f'xmlns:m="{M_NS}" xmlns:a="{A_NS}">'
                    f'<m:oMath>{run_text}</m:oMath>'
                    f'</a14:m>'
                )
                a14_elem = ET.fromstring(a14_xml)
                for r_elem in a14_elem.iter(qn('m:r')):
                    arPr = r_elem.find(qn('a:rPr'))
                    if arPr is None:
                        arPr = ET.Element(qn('a:rPr'))
                        r_elem.insert(0, arPr)
                    arPr.set('sz', str(sz_centi))
                    if arPr.get('lang') is None:
                        arPr.set('lang', 'en-US')
                    for sf in arPr.findall(qn('a:solidFill')):
                        arPr.remove(sf)
                    # fill must come BEFORE a:latin (schema order), else
                    # PowerPoint ignores the color
                    sf = arPr.makeelement(qn('a:solidFill'), {})
                    srgb = ET.SubElement(sf, qn('a:srgbClr'))
                    srgb.set('val', clr_hex)
                    arPr.insert(0, sf)
                p._p.append(a14_elem)
                continue

            run = p.add_run()
            run.text = run_text
            run.font.name = run_opts.get('font_name', 'Calibri')

            # Size: run-opt > para-opt > inherit (for inherits_from_lst
            # cases) or fall back to level default (for explicit-styled
            # paragraphs).
            run_sz = run_opts.get('size')
            sz = run_sz if run_sz is not None else para_size_override
            if sz is None and not inherits_from_lst:
                sz = size if level == 0 else sub_size
            if sz is not None:
                run.font.size = Pt(sz)

            run_clr = run_opts.get('color')
            clr = run_clr if run_clr is not None else para_color_override
            if clr is None and not inherits_from_lst:
                clr = NAVY if level == 0 else GRAY
            if clr is not None:
                run.font.color.rgb = clr

            # Bold: only set if explicitly requested.  Default = inherit
            # (lstStyle defRPr has b="0").
            run_b = run_opts.get('bold')
            b = run_b if run_b is not None else para_bold_override
            if b is not None:
                run.font.bold = b

            it = run_opts.get('italic', para_italic_override)
            if it is not None:
                run.font.italic = it
            if run_opts.get('underline'):
                run.font.underline = True
            if run_opts.get('wingdings'):
                # Add <a:sym typeface="Wingdings"/> so private-use-area
                # characters render as their Wingdings glyphs.
                rPr = run._r.find(qn('a:rPr'))
                if rPr is None:
                    rPr = run._r.makeelement(qn('a:rPr'), {})
                    run._r.insert(0, rPr)
                sym = ET.SubElement(rPr, qn('a:sym'))
                sym.set('typeface', 'Wingdings')

    return box


# --------------------------------------------------------------------------
# Layout-3 variant — diagram canvas (action title + free-form shapes below).
# Used for the agenda flowchart and the Big Picture diagram so they live
# inside the same visual chrome as content slides but render a diagram
# instead of bullets.
# --------------------------------------------------------------------------

def make_diagram_slide(prs, page_num, section_tag, title, draw_diagram):
    """Action-title slide with free-form diagram region below.

    draw_diagram: callable(slide) that renders the diagram in the body
    region (approximately y = 1.85 to 6.95).
    """
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, section_tag)
    _draw_action_title(slide, title)
    draw_diagram(slide)
    _draw_footer(slide, FOOTER_TEXT, page_num)
    return slide


# --------------------------------------------------------------------------
# Layout 5 — Poll slide (A./B./C./D. options + POLL pill, no QR box)
# --------------------------------------------------------------------------

def _draw_poll_pill(slide, *, position='top-right',
                     fill=NAVY, text_color=WHITE, dot_color=GOLD,
                     width=None, height=None, text_size=14,
                     shadow=False):
    """Small 'POLL' pill, right-aligned.

    position: 'top-right' (default, just below the top bar) or
              'bottom-right' (bottom band aligned with discussion break).
    width / height: override the default 1.05" × 0.42" pill size.
    text_size: "POLL" font size (pt).  Scales up when the pill is enlarged.
    dot_color: pass None to skip the leading accent dot.
    shadow: add a soft drop shadow to the pill (matches discussion-break
        chrome).  Default False to preserve the original flat top-right
        pill look.
    """
    pill_w = width if width is not None else Inches(1.05)
    pill_h = height if height is not None else Inches(0.42)
    pill_x = SLIDE_W - MARGIN - pill_w
    if position == 'bottom-right':
        # Bottom-aligned with where _add_discussion_break ends
        # (top=6.25, height=0.72 → bottom=6.97).  Pill height varies, so
        # pick top so the pill BOTTOM lands on the same 6.97 baseline.
        pill_y = Inches(6.97) - pill_h
    else:
        pill_y = Inches(0.55)

    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  pill_x, pill_y, pill_w, pill_h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    if shadow:
        _add_drop_shadow(shp)
    else:
        shp.shadow.inherit = False
    try:
        shp.adjustments[0] = 0.5
    except Exception:
        pass
    _add_text(slide, pill_x, pill_y, pill_w, pill_h,
              "POLL", size=text_size, bold=True, color=text_color, font="Calibri",
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if dot_color is not None:
        dot_size = Inches(0.16)
        dot_x = pill_x - Inches(0.16) - dot_size
        dot_y = pill_y + (pill_h - dot_size) // 2
        _add_rect(slide, dot_x, dot_y, dot_size, dot_size, dot_color)


def make_poll_slide(prs, page_num, section_tag, title, options, *,
                    instructions="Respond at PollEv.com/nvoigtlaender",
                    size=30, line_spacing_pts=22):
    """Layout 5 – Poll slide.

    options: list of strings (A./B./C./D. auto-numbering applied).
    """
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, section_tag)
    _draw_action_title(slide, title)
    _draw_poll_pill(slide)

    _add_bulleted_list(
        slide,
        left=MARGIN,
        top=Inches(2.0),
        width=RULE_W,
        height=Inches(4.4),
        items=options,
        size=size, color=NAVY, bullet_color=NAVY,
        line_spacing_pts=line_spacing_pts,
        autonum_scheme='alphaUcPeriod',
    )

    _add_text(slide, MARGIN, Inches(6.55), RULE_W, Inches(0.4),
              instructions, size=16, italic=True, color=GRAY,
              align=PP_ALIGN.RIGHT)

    _draw_footer(slide, FOOTER_TEXT, page_num)
    return slide


# --------------------------------------------------------------------------
# Slide-jump hyperlinks — make a run or shape clickable in slideshow
# mode so it advances PowerPoint to a different slide.  python-pptx
# doesn't expose this directly, so we manipulate the rPr / cNvPr XML
# and register the slide relationship via the public OPC API.
# --------------------------------------------------------------------------

def _add_slide_jump_hyperlink_run(source_slide, run, target_slide,
                                    *, lock_color=True, underline=True):
    """Make ``run`` a clickable hyperlink that advances to ``target_slide``.

    ``lock_color=True`` adds the Office hyperlinkcolor extension so the
    run's explicit ``solidFill`` color is preserved at render time —
    without it, PowerPoint repaints hyperlink text in the theme's
    ``hlink`` color (a light blue) regardless of what the rPr says.

    ``underline=True`` sets ``u="sng"`` so the hyperlink stays visibly
    underlined even when the color-lock suppresses the theme styling.
    """
    rId = source_slide.part.relate_to(target_slide.part, RT.SLIDE)
    rPr = run._r.get_or_add_rPr()
    if underline:
        rPr.set('u', 'sng')
    for hl in rPr.findall(qn('a:hlinkClick')):
        rPr.remove(hl)
    hlinkClick = ET.SubElement(rPr, qn('a:hlinkClick'))
    hlinkClick.set(qn('r:id'), rId)
    hlinkClick.set('action', 'ppaction://hlinksldjump')
    if lock_color:
        AHYP_NS = 'http://schemas.microsoft.com/office/drawing/2018/hyperlinkcolor'
        ext_xml = (
            f'<a:extLst xmlns:a="{A_NS}">'
            f'<a:ext uri="{{A12FA001-AC4F-418D-AE19-62706E023703}}">'
            f'<ahyp:hlinkClr xmlns:ahyp="{AHYP_NS}" val="tx"/>'
            f'</a:ext>'
            f'</a:extLst>'
        )
        hlinkClick.append(ET.fromstring(ext_xml))


def _add_slide_jump_hyperlink_shape(source_slide, shape, target_slide):
    """Make ``shape`` clickable so the whole shape jumps to ``target_slide``."""
    rId = source_slide.part.relate_to(target_slide.part, RT.SLIDE)
    nvSpPr = shape._element.find(qn('p:nvSpPr'))
    cNvPr = nvSpPr.find(qn('p:cNvPr')) if nvSpPr is not None else None
    if cNvPr is None:
        return
    for hl in cNvPr.findall(qn('a:hlinkClick')):
        cNvPr.remove(hl)
    hlinkClick = ET.SubElement(cNvPr, qn('a:hlinkClick'))
    hlinkClick.set(qn('r:id'), rId)
    hlinkClick.set('action', 'ppaction://hlinksldjump')


# --------------------------------------------------------------------------
# Image helpers — embed source-deck images into the new deck.
# Images are pre-extracted to _source_images/slide{N}_{rId}.{ext}.
# --------------------------------------------------------------------------

SRC_IMG_DIR = Path(__file__).parent / "_source_images"


def _add_drop_shadow(shape, *, blur="50800", dist="38100",
                      direction="2700000", alpha="45000"):
    """Add a soft drop shadow to any shape that exposes spPr (pictures,
    rectangles, rounded rects).  Default: 4pt blur, 3pt offset, 45° down-
    right, 45% opacity black.  Used deck-wide for figures/boxes."""
    try:
        spPr = shape._element.spPr
    except AttributeError:
        # Fallback for shapes whose XML element exposes spPr only via find()
        spPr = shape._element.find(qn('p:spPr'))
    if spPr is None:
        return shape
    for old in spPr.findall(qn('a:effectLst')):
        spPr.remove(old)
    effLst = ET.SubElement(spPr, qn('a:effectLst'))
    outerShdw = ET.SubElement(effLst, qn('a:outerShdw'))
    outerShdw.set('blurRad', str(blur))
    outerShdw.set('dist', str(dist))
    outerShdw.set('dir', str(direction))
    outerShdw.set('algn', 'tl')
    outerShdw.set('rotWithShape', '0')
    rgb = ET.SubElement(outerShdw, qn('a:srgbClr'))
    rgb.set('val', '000000')
    a = ET.SubElement(rgb, qn('a:alpha'))
    a.set('val', str(alpha))
    return shape


def _add_graphicframe_shadow(slide, left, top, width, height, *,
                              shadow_alpha=45000):
    """White backing rectangle with an outerShdw effect, behind a
    graphicFrame (table or chart).  graphicFrames can't host
    a:effectLst directly; this rect supplies the shadow projected
    OUTSIDE its bounds.

    Charts have transparent plot areas by default, so a coloured backing
    bleeds through and tints the figure.  Use white instead — the chart
    sees white through its own transparent areas (clean background) and
    the shadow renders only at the visible edges.  Call BEFORE adding
    the table/chart so z-order is correct.
    """
    shdw = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        int(left), int(top), int(width), int(height),
    )
    shdw.fill.solid()
    shdw.fill.fore_color.rgb = WHITE
    shdw.line.fill.background()
    shdw.shadow.inherit = False
    sp_pr = shdw._element.spPr
    # Strip any default <a:effectLst/> python-pptx may have inserted
    # (duplicate effectLst makes PowerPoint refuse to open the file).
    for old in sp_pr.findall(qn('a:effectLst')):
        sp_pr.remove(old)
    effLst = ET.SubElement(sp_pr, qn('a:effectLst'))
    outerShdw = ET.SubElement(effLst, qn('a:outerShdw'))
    outerShdw.set('blurRad', '50800')
    outerShdw.set('dist', '38100')
    outerShdw.set('dir', '2700000')
    outerShdw.set('algn', 'tl')
    outerShdw.set('rotWithShape', '0')
    rgb = ET.SubElement(outerShdw, qn('a:srgbClr'))
    rgb.set('val', '000000')
    a = ET.SubElement(rgb, qn('a:alpha'))
    a.set('val', str(int(shadow_alpha)))
    return shdw


def _add_source_image(slide, src_slide_no, rid, *, left, top, width=None,
                      height=None, shadow=True, rounded=False):
    """Place a source-deck image on the new slide.

    `shadow=True` (default) adds a soft drop shadow so figures pop off the
    background — applied deck-wide per the latest visual direction.  Set
    `shadow=False` for niche cases (transparent PNGs, screenshots that
    already include a shadow, etc.).

    `rounded=True` gives real photographs the deck-standard rounded corners
    (routes through :func:`_apply_picture_style`, which sets both the
    rounded geometry and the drop shadow).  Leave False for logos,
    screenshots, and framed images per the Pictures guideline.
    """
    candidates = list(SRC_IMG_DIR.glob(f"slide{src_slide_no}_{rid}.*"))
    if not candidates:
        return None
    img = candidates[0]
    kwargs = {"left": int(left), "top": int(top)}
    if width is not None:
        kwargs["width"] = int(width)
    if height is not None:
        kwargs["height"] = int(height)
    pic = slide.shapes.add_picture(str(img), **kwargs)
    if rounded:
        _apply_picture_style(pic, corner_pct=6)
    elif shadow:
        _add_drop_shadow(pic)
    return pic


def _apply_picture_style(pic, *, corner_pct=8,
                          shadow_blur=50800, shadow_dist=38100,
                          shadow_dir=2700000, shadow_alpha=50000):
    """Apply rounded corners + drop shadow to a picture shape.

    corner_pct: rounded-corner radius as percent of shorter side (8 ≈ subtle).
    shadow_blur/dist: EMU; defaults give a soft 4pt blur, 3pt offset.
    shadow_dir: 2700000 = 45° down-right (standard).
    shadow_alpha: 50000 = 50% opacity black shadow.
    """
    spPr = pic._element.find(qn('p:spPr'))
    if spPr is None:
        return pic
    # Replace any existing prstGeom with roundRect at corner_pct
    for old in spPr.findall(qn('a:prstGeom')):
        spPr.remove(old)
    prstGeom = ET.Element(qn('a:prstGeom'))
    prstGeom.set('prst', 'roundRect')
    avLst = ET.SubElement(prstGeom, qn('a:avLst'))
    gd = ET.SubElement(avLst, qn('a:gd'))
    gd.set('name', 'adj')
    gd.set('fmla', f'val {int(corner_pct * 1000)}')
    # prstGeom must come after a:xfrm
    xfrm = spPr.find(qn('a:xfrm'))
    if xfrm is not None:
        xfrm.addnext(prstGeom)
    else:
        spPr.insert(0, prstGeom)
    # Replace any existing effectLst with a single outer shadow
    for old in spPr.findall(qn('a:effectLst')):
        spPr.remove(old)
    effectLst = ET.SubElement(spPr, qn('a:effectLst'))
    outerShdw = ET.SubElement(effectLst, qn('a:outerShdw'))
    outerShdw.set('blurRad', str(int(shadow_blur)))
    outerShdw.set('dist', str(int(shadow_dist)))
    outerShdw.set('dir', str(int(shadow_dir)))
    outerShdw.set('algn', 'tl')
    outerShdw.set('rotWithShape', '0')
    rgb = ET.SubElement(outerShdw, qn('a:srgbClr'))
    rgb.set('val', '000000')
    alpha = ET.SubElement(rgb, qn('a:alpha'))
    alpha.set('val', str(int(shadow_alpha)))
    return pic


# --------------------------------------------------------------------------
# Illustrative callout boxes — the "major-concept" badges, "Teaching Note"
# bars, and bottom-of-slide takeaway bands that recur throughout the source
# deck.  Replicating them faithfully (in template colors) is what makes the
# slides feel like the original, not a stripped-down rewrite.
# --------------------------------------------------------------------------

def _add_takeaway_bar(slide, text, *, top=Inches(6.4), width=None,
                       height=Inches(0.55), left=None,
                       fill=GOLD, text_color=WHITE,
                       size=20, font="Calibri", bold=True,
                       rounded=False, shadow=False):
    """Bottom-of-slide takeaway band — the 'major concept' callout.

    rounded=True  → ROUNDED_RECTANGLE with ~30% corner radius.
    shadow=True   → soft drop shadow (off by default to keep the flat
                    look on the majority of slides).
    left=None     → horizontally centered. Pass an EMU/Inches value to
                    pin the bar's left edge (used for hand-positioned
                    takeaways).
    """
    if width is None:
        width = Inches(9.6)
    if left is None:
        left = (SLIDE_W - width) // 2
    if not (rounded or shadow):
        return _add_filled_box(slide, left, top, width, height, text,
                                fill=fill, text_color=text_color,
                                size=size, bold=bold, font=font)
    # Inline build for the rounded / shadowed variant
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, int(left), int(top), int(width), int(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    if rounded:
        try: shape.adjustments[0] = 0.30
        except Exception: pass
    shape.shadow.inherit = False
    if shadow:
        _add_drop_shadow(shape)
    tf = shape.text_frame
    tf.margin_left = tf.margin_right = Inches(0.1)
    tf.margin_top = tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = text_color
    return shape


def _add_teaching_note(slide, text, *, top=Inches(6.6), width=None,
                        height=Inches(0.6), left=None,
                        rounded=False, pdf_icon=False, label_color=GOLD,
                        fill_rgb=None):
    """External-document reference card.

    Visually distinct from in-slide callouts: cream/parchment fill, navy
    DASHED border, italic navy text, with a small page-icon glyph on the
    left — clearly signals 'this links to an external Teaching Note doc'.

    Optional styling for the slide-32 (page 33) treatment:
      • ``left``       — pin to an absolute X (default: horizontally center)
      • ``rounded``    — rounded corners + soft drop shadow (default flat)
      • ``pdf_icon``   — render the page glyph as a white folded-corner
                        with bold "PDF" text, sized to nearly fill the card
                        (default: small navy folded-corner, no text)
      • ``label_color``— color of the "SEE TEACHING NOTE →" prefix
                        (default GOLD; pass NAVY for the page-33 look).
      • ``fill_rgb``   — override the cream/parchment card fill
                        (default RGBColor(0xF4, 0xF1, 0xEA)).
    """
    if width is None:
        width = Inches(8.0)
    if left is None:
        left = int((SLIDE_W - width) // 2)
    left, top, width, height = int(left), int(top), int(width), int(height)

    # Card: cream fill, dashed navy border.  ``rounded=True`` switches the
    # base shape to ROUNDED_RECTANGLE and adds a soft drop shadow.
    base_shape = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    card = slide.shapes.add_shape(base_shape, left, top, width, height)
    if rounded:
        try: card.adjustments[0] = 0.20
        except Exception: pass
    card.fill.solid()
    card.fill.fore_color.rgb = (fill_rgb if fill_rgb is not None
                                  else RGBColor(0xF4, 0xF1, 0xEA))  # parchment cream default
    card.line.color.rgb = NAVY
    card.line.width = Pt(1.25)
    card.shadow.inherit = False
    if rounded:
        _add_drop_shadow(card)
    # Dashed line style
    ln = card.line._get_or_add_ln()
    # Remove any existing prstDash
    for el in ln.findall(qn('a:prstDash')):
        ln.remove(el)
    prstDash = ET.SubElement(ln, qn('a:prstDash'))
    prstDash.set('val', 'dash')

    # Empty text on the card; we add the icon + text as separate textboxes
    tf = card.text_frame
    tf.text = ""

    # Small "page" icon on the left – a folded-corner shape.  In default
    # mode it's a small filled-navy glyph (logo-ish); in pdf_icon mode it's
    # a larger WHITE sheet with thin navy border + bold "PDF" text so it
    # reads as a generic PDF document at a glance.
    icon_size = Inches(0.5) if pdf_icon else Inches(0.4)
    icon_x = left + Inches(0.2)
    icon_y = top + (height - icon_size) // 2
    page = slide.shapes.add_shape(MSO_SHAPE.FOLDED_CORNER,
                                   int(icon_x), int(icon_y),
                                   int(icon_size), int(icon_size))
    page.fill.solid()
    if pdf_icon:
        page.fill.fore_color.rgb = WHITE
        page.line.color.rgb = NAVY
        page.line.width = Pt(0.75)
    else:
        page.fill.fore_color.rgb = NAVY
        page.line.fill.background()
    page.shadow.inherit = False
    if pdf_icon:
        # "PDF" label inside the white sheet.
        ptf = page.text_frame
        ptf.word_wrap = False
        ptf.margin_left = 0
        ptf.margin_right = 0
        ptf.margin_top = 0
        ptf.margin_bottom = 0
        ptf.vertical_anchor = MSO_ANCHOR.MIDDLE
        pp = ptf.paragraphs[0]
        pp.alignment = PP_ALIGN.CENTER
        pr = pp.add_run()
        pr.text = "PDF"
        pr.font.name = "Calibri"
        pr.font.size = Pt(11)
        pr.font.bold = True
        pr.font.color.rgb = NAVY

    # Label text (italic, navy) inside the card to the right of the icon
    txt_left = int(icon_x + icon_size + Inches(0.2))
    txt_w = int(width - (icon_x + icon_size + Inches(0.4) - left))
    label_box = slide.shapes.add_textbox(txt_left, top,
                                          txt_w, height)
    label_tf = label_box.text_frame
    label_tf.word_wrap = True
    label_tf.margin_left = 0
    label_tf.margin_right = 0
    label_tf.margin_top = 0
    label_tf.margin_bottom = 0
    label_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = label_tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    # First run: "See teaching note:" prefix in caller-supplied color
    # (GOLD by default; NAVY for the slide-33 treatment).
    r1 = p.add_run()
    r1.text = "SEE TEACHING NOTE  →  "
    r1.font.name = "Calibri"
    r1.font.size = Pt(12)
    r1.font.bold = True
    r1.font.color.rgb = label_color
    # Second run: the title of the note, italic navy
    r2 = p.add_run()
    r2.text = text
    r2.font.name = "Calibri"
    r2.font.size = Pt(16)
    r2.font.italic = True
    r2.font.bold = True
    r2.font.color.rgb = NAVY
    return card


def _add_discussion_break(slide, *, top=Inches(6.25), width=Inches(4.8),
                           left=None, text="Discussion Break"):
    """Rounded-parallelogram 'discussion break' badge (bottom-right).

    Custom-geometry shape: top and bottom edges are horizontal; the left
    and right edges slant at 45° in real space (skew = height of the
    shape).  All four corners are slightly rounded.  Gold fill, navy
    bold text, soft drop shadow.

    left=None → pin to the right edge with the default MARGIN. Pass an
    EMU/Inches value to override (used for hand-positioned badges).
    """
    height = Inches(0.72)
    if left is None:
        left = SLIDE_W - MARGIN - width
    left, top, width, height = int(left), int(top), int(width), int(height)
    # Compute skew in path-coordinate units so that left/right sides slant
    # at 45° in REAL space:  horizontal offset of top edge == shape height.
    skew = int(100000 * height / width) if width else 15000
    skew = min(max(skew, 6000), 35000)
    r = 5000          # corner radius in path units — gentle rounding

    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = GOLD
    shp.line.fill.background()
    shp.shadow.inherit = False
    # Strip the default prstGeom – we'll inject custGeom instead.
    spPr = shp._element.spPr
    for old in spPr.findall(qn('a:prstGeom')):
        spPr.remove(old)
    rs = int(r * skew / 100000)
    # IMPORTANT: <a:rect> defines the TEXT bounding rectangle inside the
    # custom geometry.  Set it to the parallelogram's inscribed rectangle
    # (from TL vertex to BR vertex) so PowerPoint won't render text past
    # the slanted edges, regardless of the text frame's own lIns/rIns.
    custgeom_xml = (
        f'<a:custGeom xmlns:a="{A_NS}">'
        f'<a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/>'
        f'<a:rect l="{skew}" t="0" r="{100000-skew}" b="100000"/>'
        f'<a:pathLst><a:path w="100000" h="100000">'
        # Top-left rounded corner (vertex at (skew, 0))
        f'<a:moveTo><a:pt x="{skew+r}" y="0"/></a:moveTo>'
        # Top edge
        f'<a:lnTo><a:pt x="{100000-r}" y="0"/></a:lnTo>'
        # Top-right corner round
        f'<a:cubicBezTo>'
        f'<a:pt x="100000" y="0"/><a:pt x="100000" y="0"/>'
        f'<a:pt x="{100000-rs}" y="{r}"/>'
        f'</a:cubicBezTo>'
        # Right slanted side (down-left)
        f'<a:lnTo><a:pt x="{100000-skew+rs}" y="{100000-r}"/></a:lnTo>'
        # Bottom-right corner round (vertex at (100000-skew, 100000))
        f'<a:cubicBezTo>'
        f'<a:pt x="{100000-skew}" y="100000"/>'
        f'<a:pt x="{100000-skew}" y="100000"/>'
        f'<a:pt x="{100000-skew-r}" y="100000"/>'
        f'</a:cubicBezTo>'
        # Bottom edge
        f'<a:lnTo><a:pt x="{r}" y="100000"/></a:lnTo>'
        # Bottom-left corner round (vertex at (0, 100000))
        f'<a:cubicBezTo>'
        f'<a:pt x="0" y="100000"/><a:pt x="0" y="100000"/>'
        f'<a:pt x="{rs}" y="{100000-r}"/>'
        f'</a:cubicBezTo>'
        # Left slanted side (up-right)
        f'<a:lnTo><a:pt x="{skew-rs}" y="{r}"/></a:lnTo>'
        # Top-left corner round (close the path)
        f'<a:cubicBezTo>'
        f'<a:pt x="{skew}" y="0"/><a:pt x="{skew}" y="0"/>'
        f'<a:pt x="{skew+r}" y="0"/>'
        f'</a:cubicBezTo>'
        f'<a:close/>'
        f'</a:path></a:pathLst>'
        f'</a:custGeom>'
    )
    custgeom = ET.fromstring(custgeom_xml)
    # Insert custGeom right after a:xfrm (schema order)
    xfrm = spPr.find(qn('a:xfrm'))
    if xfrm is not None:
        xfrm.addnext(custgeom)
    else:
        spPr.insert(0, custgeom)

    # Drop shadow (45° down-right, 50% opacity).
    for old in spPr.findall(qn('a:effectLst')):
        spPr.remove(old)
    effectLst = ET.SubElement(spPr, qn('a:effectLst'))
    outerShdw = ET.SubElement(effectLst, qn('a:outerShdw'))
    outerShdw.set('blurRad', '50800')
    outerShdw.set('dist', '38100')
    outerShdw.set('dir', '2700000')
    outerShdw.set('algn', 'tl')
    outerShdw.set('rotWithShape', '0')
    rgb = ET.SubElement(outerShdw, qn('a:srgbClr'))
    rgb.set('val', '000000')
    alpha = ET.SubElement(rgb, qn('a:alpha'))
    alpha.set('val', '50000')

    # Text — render as a SEPARATE textbox overlaid on top of the
    # parallelogram, positioned exactly inside the inscribed rectangle.
    # This decouples text placement from the shape geometry and avoids
    # PowerPoint rendering the run past the slanted edges (which can
    # happen with the in-shape text frame on some PowerPoint versions
    # even with <a:rect> set inside custGeom).
    # In real EMU, the skew equals the shape height (45° slant), so the
    # inscribed rectangle spans (left + height) → (left + width - height).
    skew_emu = height
    ins_left = left + skew_emu
    ins_top = top
    ins_w = width - 2 * skew_emu
    ins_h = height
    txt = slide.shapes.add_textbox(int(ins_left), int(ins_top),
                                     int(ins_w), int(ins_h))
    tf = txt.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(28)         # bumped 20 → 28 on 2026-05-16
    run.font.bold = True
    run.font.color.rgb = NAVY
    return shp


def _add_callout_box(slide, left, top, width, height, text, *,
                      fill=GOLD, text_color=WHITE, size=14, bold=True,
                      rounded=False, shadow=False):
    """Small free-form annotation/callout (e.g., 'plot the slope', 'Revenue
    per car net of material cost').  Used to mark a graph or sub-region.

    ``rounded``/``shadow`` (default off) route to the rounded + drop-shadow
    variant for the deck-wide box treatment."""
    if rounded or shadow:
        return _add_rounded_filled_box(
            slide, left, top, width, height, text,
            fill=fill, text_color=text_color,
            size=size, bold=bold, font="Calibri",
            corner_pct=0.12, shadow=shadow)
    return _add_filled_box(slide, left, top, width, height, text,
                            fill=fill, text_color=text_color,
                            size=size, bold=bold, font="Calibri")


def _add_anchor_burst(slide, left, top, width, height,
                       top_text, bottom_text=None, extra_text=None,
                       *, fill=GOLD, text_color=NAVY,
                       top_size=14, bottom_size=10):
    """12-point star background + a separate text-box overlay.

    The star is purely decorative; the text lives in a normal
    rectangular text box layered on top so the text isn't constrained
    by the star's geometry (no risk of bleeding into the points).
    Reusable on every slide where MB = MC is being invoked, so the same
    visual pattern carries over.
    """
    left, top, width, height = int(left), int(top), int(width), int(height)

    # 1. Decorative star background (no text)
    shp = slide.shapes.add_shape(MSO_SHAPE.STAR_12_POINT, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = NAVY
    shp.line.width = Pt(1.0)
    shp.shadow.inherit = False
    # Soft drop shadow — added 2026-05-15 per user request, so the MB=MC
    # star reads as lifted off the slide like every other content shape.
    _add_drop_shadow(shp)
    # Suppress any auto-inserted text frame contents on the shape itself.
    shp.text_frame.text = ""

    # 2. Overlay text box, sized to the star's inscribed body
    inner_w = int(width * 0.65)
    inner_h = int(height * 0.55)
    inner_x = left + (width - inner_w) // 2
    inner_y = top + (height - inner_h) // 2

    box = slide.shapes.add_textbox(inner_x, inner_y, inner_w, inner_h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = top_text
    r1.font.name = 'Calibri'
    r1.font.size = Pt(top_size)
    r1.font.bold = True
    r1.font.color.rgb = text_color
    if bottom_text:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = bottom_text
        r2.font.name = 'Calibri'
        r2.font.size = Pt(bottom_size)
        r2.font.italic = True
        r2.font.bold = True
        r2.font.color.rgb = text_color
    if extra_text:
        p3 = tf.add_paragraph()
        p3.alignment = PP_ALIGN.CENTER
        r3 = p3.add_run()
        r3.text = extra_text
        r3.font.name = 'Calibri'
        r3.font.size = Pt(bottom_size)
        r3.font.italic = True
        r3.font.bold = True
        r3.font.color.rgb = text_color
    return shp


# --------------------------------------------------------------------------
# OMML (Office Math Markup Language) equation helper – gives formulas a
# proper TeX-style render with italic variables, stacked fractions, real
# subscripts/superscripts.  Uses Cambria Math (the standard PPT math font).
# --------------------------------------------------------------------------

M_NS = 'http://schemas.openxmlformats.org/officeDocument/2006/math'
A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
A14_NS = 'http://schemas.microsoft.com/office/drawing/2010/main'


def _omml_fill(color):
    """Build the inner ``<a:solidFill>`` clause for an OMML run's rPr.

    ``color`` may be an ``RGBColor`` instance (or a 3-tuple of ints).
    Returns an empty string when ``color`` is None so callers can splice
    the result into rPr unconditionally.
    """
    if color is None:
        return ''
    return (
        f'<a:solidFill><a:srgbClr val="'
        f'{color[0]:02X}{color[1]:02X}{color[2]:02X}'
        f'"/></a:solidFill>'
    )


def _omml_run(text, *, color=None):
    """OMML run for an italic variable (default math style).

    Inside an oMath, italic style is the math default for Latin letters;
    we leave m:rPr out entirely so the Cambria Math italic comes through.
    The a:rPr applies drawing-level font sizing/coloring.  Pass ``color``
    to tint the run (e.g., green ΔL / ΔQ in the slide-14 Convention box).
    """
    return (
        f'<m:r xmlns:m="{M_NS}">'
        f'<a:rPr xmlns:a="{A_NS}" lang="en-US" b="0" i="1">'
        f'{_omml_fill(color)}'
        f'<a:latin typeface="Cambria Math"/>'
        f'<a:ea typeface="Cambria Math"/>'
        f'</a:rPr>'
        f'<m:t>{text}</m:t>'
        f'</m:r>'
    )


def _omml_text(text, *, color=None):
    """Upright-style OMML run (for operators, numbers, acronyms).

    Force plain (upright) style via <m:rPr><m:sty m:val="p"/></m:rPr> – this
    is the documented way to disable the math-default italics for the
    enclosed run.  Pass ``color`` to tint the run.
    """
    return (
        f'<m:r xmlns:m="{M_NS}">'
        f'<m:rPr><m:sty m:val="p"/></m:rPr>'
        f'<a:rPr xmlns:a="{A_NS}" lang="en-US" b="0" i="0">'
        f'{_omml_fill(color)}'
        f'<a:latin typeface="Cambria Math"/>'
        f'</a:rPr>'
        f'<m:t xml:space="preserve">{text}</m:t>'
        f'</m:r>'
    )


def _omml_sub(base, sub):
    """OMML subscript: base with subscript expression."""
    return (
        f'<m:sSub xmlns:m="{M_NS}">'
        f'<m:sSubPr><m:ctrlPr>'
        f'<a:rPr xmlns:a="{A_NS}" lang="en-US" i="1">'
        f'<a:latin typeface="Cambria Math"/></a:rPr>'
        f'</m:ctrlPr></m:sSubPr>'
        f'<m:e>{base}</m:e>'
        f'<m:sub>{sub}</m:sub>'
        f'</m:sSub>'
    )


def _omml_frac(num, den):
    """OMML stacked fraction: num / den."""
    return (
        f'<m:f xmlns:m="{M_NS}">'
        f'<m:fPr><m:ctrlPr>'
        f'<a:rPr xmlns:a="{A_NS}" lang="en-US" i="1">'
        f'<a:latin typeface="Cambria Math"/></a:rPr>'
        f'</m:ctrlPr></m:fPr>'
        f'<m:num>{num}</m:num>'
        f'<m:den>{den}</m:den>'
        f'</m:f>'
    )


def _omml_sup(base, sup):
    """OMML superscript: base^sup (e.g. Q²)."""
    return (
        f'<m:sSup xmlns:m="{M_NS}">'
        f'<m:sSupPr><m:ctrlPr>'
        f'<a:rPr xmlns:a="{A_NS}" lang="en-US" i="1">'
        f'<a:latin typeface="Cambria Math"/></a:rPr>'
        f'</m:ctrlPr></m:sSupPr>'
        f'<m:e>{base}</m:e>'
        f'<m:sup>{sup}</m:sup>'
        f'</m:sSup>'
    )


def _add_dashed_gridlines(axis_el):
    """Dashed light-grey major gridlines on an axis (Cobb-Douglas style)."""
    for old in axis_el.findall(qn('c:majorGridlines')):
        axis_el.remove(old)
    gl = ET.Element(qn('c:majorGridlines'))
    sp = ET.SubElement(gl, qn('c:spPr'))
    ln = ET.SubElement(sp, qn('a:ln'))
    ln.set('w', '9525')
    ln.set('cap', 'flat'); ln.set('cmpd', 'sng'); ln.set('algn', 'ctr')
    fill = ET.SubElement(ln, qn('a:solidFill'))
    clr = ET.SubElement(fill, qn('a:srgbClr')); clr.set('val', 'C8CDD3')
    dash = ET.SubElement(ln, qn('a:prstDash')); dash.set('val', 'dash')
    axpos = axis_el.find(qn('c:axPos'))
    if axpos is not None:
        axpos.addnext(gl)


def _align_x_labels_with_ticks(value_axis):
    """Set <c:crossBetween val="midCat"/> on a value axis so the category
    labels and tick marks align (default for line charts in OOXML is
    "between", which leaves labels visually between adjacent ticks).
    """
    val_el = value_axis._element
    for old in val_el.findall(qn('c:crossBetween')):
        val_el.remove(old)
    cb = ET.Element(qn('c:crossBetween'))
    cb.set('val', 'midCat')
    # Schema position: c:crossBetween follows c:crosses(At) and precedes
    # c:majorUnit.  Insert before c:majorUnit if present; else append.
    mu_el = val_el.find(qn('c:majorUnit'))
    if mu_el is not None:
        mu_el.addprevious(cb)
    else:
        val_el.append(cb)


def _make_simple_line_chart(slide, x, y, w, h, categories, values, *,
                              line_color, x_title, y_title,
                              y_min=0, y_max=None, y_unit=None,
                              marker='circle'):
    """Single-series line+markers chart with dashed light-grey gridlines.

    Same visual conventions as slide 11 (Calibri navy labels, dashed C8CDD3
    major gridlines, marker size 7) but no legend/title.
    """
    cd = CategoryChartData()
    cd.categories = list(categories)
    cd.add_series("Y", values)
    # Drop-shadow rectangle behind the chart (graphicFrames can't host shadow).
    _add_graphicframe_shadow(slide, x, y, w, h)
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE,
        int(x), int(y), int(w), int(h), cd,
    )
    chart = chart_shape.chart
    chart.has_title = False
    chart.has_legend = False

    clr_hex = f'{line_color[0]:02X}{line_color[1]:02X}{line_color[2]:02X}'

    for series in chart.series:
        line = series.format.line
        line.color.rgb = line_color
        line.width = Pt(2.5)
        ser_xml = series._element
        for old in ser_xml.findall(qn('c:marker')):
            ser_xml.remove(old)
        m = ET.SubElement(ser_xml, qn('c:marker'))
        sym = ET.SubElement(m, qn('c:symbol')); sym.set('val', marker)
        sz_el = ET.SubElement(m, qn('c:size')); sz_el.set('val', '7')
        sp = ET.SubElement(m, qn('c:spPr'))
        fl = ET.SubElement(sp, qn('a:solidFill'))
        rg = ET.SubElement(fl, qn('a:srgbClr')); rg.set('val', clr_hex)
        ln = ET.SubElement(sp, qn('a:ln'))
        lf = ET.SubElement(ln, qn('a:solidFill'))
        lr = ET.SubElement(lf, qn('a:srgbClr')); lr.set('val', clr_hex)
        # disable smoothing (straight segments between points)
        for sm in ser_xml.findall(qn('c:smooth')):
            ser_xml.remove(sm)
        smooth = ET.SubElement(ser_xml, qn('c:smooth'))
        smooth.set('val', '0')

    # Axes – axis titles in BOLD ITALIC navy (per course CLAUDE.md);
    # tick labels in regular Calibri navy.
    cat = chart.category_axis
    cat.tick_labels.font.name = "Calibri"
    cat.tick_labels.font.size = Pt(10)
    cat.tick_labels.font.color.rgb = NAVY
    cat.has_title = True
    cat.axis_title.text_frame.text = x_title
    ar = cat.axis_title.text_frame.paragraphs[0].runs[0]
    ar.font.name = "Calibri"; ar.font.size = Pt(12)
    ar.font.bold = True; ar.font.italic = True
    ar.font.color.rgb = NAVY

    val = chart.value_axis
    val.tick_labels.font.name = "Calibri"
    val.tick_labels.font.size = Pt(10)
    val.tick_labels.font.color.rgb = NAVY
    val.has_title = True
    val.axis_title.text_frame.text = y_title
    ar = val.axis_title.text_frame.paragraphs[0].runs[0]
    ar.font.name = "Calibri"; ar.font.size = Pt(12)
    ar.font.bold = True; ar.font.italic = True
    ar.font.color.rgb = NAVY
    if y_max is not None:
        val.minimum_scale = y_min
        val.maximum_scale = y_max
    if y_unit is not None:
        val.major_unit = y_unit

    # Align X-axis category labels with tick marks (default OOXML places
    # them in the gaps between ticks).
    _align_x_labels_with_ticks(val)

    _add_dashed_gridlines(cat._element)
    _add_dashed_gridlines(val._element)
    return chart_shape


def _make_multi_line_chart(slide, x, y, w, h, categories, series, *,
                             x_title, y_title,
                             y_min=0, y_max=None, y_unit=None,
                             legend=True, legend_pos=('0.08', '0.10', '0.20', '0.20')):
    """Multi-series line+markers chart with the deck's standard styling.

    series: list of (name, values, color: RGBColor, marker: str) tuples.
    legend_pos: (x, y, w, h) in chart-fraction units (str) – top-left default.
    """
    _add_graphicframe_shadow(slide, x, y, w, h)
    cd = CategoryChartData()
    cd.categories = list(categories)
    for name, values, _color, _marker in series:
        cd.add_series(name, values)
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE,
        int(x), int(y), int(w), int(h), cd,
    )
    chart = chart_shape.chart
    chart.has_title = False
    chart.has_legend = legend

    # Per-series styling: line color + marker
    for idx, ser in enumerate(chart.series):
        name, values, color, marker = series[idx]
        line = ser.format.line
        line.color.rgb = color
        line.width = Pt(2.5)
        ser_xml = ser._element
        clr_hex = f'{color[0]:02X}{color[1]:02X}{color[2]:02X}'
        for old in ser_xml.findall(qn('c:marker')):
            ser_xml.remove(old)
        m = ET.SubElement(ser_xml, qn('c:marker'))
        sym = ET.SubElement(m, qn('c:symbol')); sym.set('val', marker)
        sz_el = ET.SubElement(m, qn('c:size')); sz_el.set('val', '7')
        sp = ET.SubElement(m, qn('c:spPr'))
        fl = ET.SubElement(sp, qn('a:solidFill'))
        rg = ET.SubElement(fl, qn('a:srgbClr')); rg.set('val', clr_hex)
        ln = ET.SubElement(sp, qn('a:ln'))
        lf = ET.SubElement(ln, qn('a:solidFill'))
        lr = ET.SubElement(lf, qn('a:srgbClr')); lr.set('val', clr_hex)
        for sm in ser_xml.findall(qn('c:smooth')):
            ser_xml.remove(sm)
        smooth = ET.SubElement(ser_xml, qn('c:smooth'))
        smooth.set('val', '0')

    # Axes – bold italic navy titles per course style
    cat = chart.category_axis
    cat.tick_labels.font.name = "Calibri"
    cat.tick_labels.font.size = Pt(10)
    cat.tick_labels.font.color.rgb = NAVY
    cat.has_title = True
    cat.axis_title.text_frame.text = x_title
    ar = cat.axis_title.text_frame.paragraphs[0].runs[0]
    ar.font.name = "Calibri"; ar.font.size = Pt(12)
    ar.font.bold = True; ar.font.italic = True; ar.font.color.rgb = NAVY

    val = chart.value_axis
    val.tick_labels.font.name = "Calibri"
    val.tick_labels.font.size = Pt(10)
    val.tick_labels.font.color.rgb = NAVY
    val.has_title = True
    val.axis_title.text_frame.text = y_title
    ar = val.axis_title.text_frame.paragraphs[0].runs[0]
    ar.font.name = "Calibri"; ar.font.size = Pt(12)
    ar.font.bold = True; ar.font.italic = True; ar.font.color.rgb = NAVY
    if y_max is not None:
        val.minimum_scale = y_min
        val.maximum_scale = y_max
    if y_unit is not None:
        val.major_unit = y_unit

    # Align X-axis category labels with tick marks.
    _align_x_labels_with_ticks(val)

    _add_dashed_gridlines(cat._element)
    _add_dashed_gridlines(val._element)

    # Legend top-left inside plot with white fill
    if legend:
        leg_el = chart.legend._element
        chart.legend.font.name = "Calibri"
        chart.legend.font.size = Pt(11)
        chart.legend.font.color.rgb = NAVY
        chart.legend.include_in_layout = False
        # Strip default legendPos / layout, replace
        for old in leg_el.findall(qn('c:layout')):
            leg_el.remove(old)
        for old in leg_el.findall(qn('c:legendPos')):
            leg_el.remove(old)
        pos = ET.SubElement(leg_el, qn('c:legendPos')); pos.set('val', 'tr')
        leg_el.remove(pos); leg_el.insert(0, pos)
        # manualLayout positions in chart-fraction units
        layout = ET.Element(qn('c:layout'))
        ml = ET.SubElement(layout, qn('c:manualLayout'))
        ET.SubElement(ml, qn('c:xMode')).set('val', 'edge')
        ET.SubElement(ml, qn('c:yMode')).set('val', 'edge')
        ET.SubElement(ml, qn('c:x')).set('val', legend_pos[0])
        ET.SubElement(ml, qn('c:y')).set('val', legend_pos[1])
        ET.SubElement(ml, qn('c:w')).set('val', legend_pos[2])
        ET.SubElement(ml, qn('c:h')).set('val', legend_pos[3])
        pos.addnext(layout)
        # White fill behind legend
        for old in leg_el.findall(qn('c:spPr')):
            leg_el.remove(old)
        leg_spPr = ET.Element(qn('c:spPr'))
        sf = ET.SubElement(leg_spPr, qn('a:solidFill'))
        clr = ET.SubElement(sf, qn('a:srgbClr')); clr.set('val', 'FFFFFF')
        ln = ET.SubElement(leg_spPr, qn('a:ln')); ln.set('w', '6350')
        lf = ET.SubElement(ln, qn('a:solidFill'))
        lc = ET.SubElement(lf, qn('a:srgbClr')); lc.set('val', '0B2B4E')
        layout.addnext(leg_spPr)
    return chart_shape


def _make_xy_line_chart(slide, x, y, w, h, *, series, x_title, y_title,
                          x_min=0, x_max=None, x_unit=None,
                          y_min=0, y_max=None, y_unit=None,
                          legend=False, legend_pos=None, smooth=False):
    """XY-scatter line+markers chart with the deck's standard styling.

    Use when you need data points to sit at arbitrary X-positions (e.g.,
    plotting MPL at the midpoint of each L-interval) while keeping tick
    marks at standard L-values.  Both axes are numeric value axes.

    series: list of ``(name, [(x, y), ...], color: RGBColor, marker: str)``.
    smooth: True for XY_SCATTER_SMOOTH (cubic spline through points),
            False for XY_SCATTER_LINES (straight segments).
    """
    _add_graphicframe_shadow(slide, x, y, w, h)
    cd = XyChartData()
    for name, points, _color, _marker in series:
        s = cd.add_series(name)
        for px, py in points:
            s.add_data_point(px, py)
    chart_type = (XL_CHART_TYPE.XY_SCATTER_SMOOTH if smooth
                  else XL_CHART_TYPE.XY_SCATTER_LINES)
    chart_shape = slide.shapes.add_chart(
        chart_type,
        int(x), int(y), int(w), int(h), cd,
    )
    chart = chart_shape.chart
    chart.has_title = False
    chart.has_legend = legend

    # Per-series styling
    for idx, ser in enumerate(chart.series):
        name, points, color, marker = series[idx]
        line = ser.format.line
        line.color.rgb = color
        line.width = Pt(2.5)
        ser_xml = ser._element
        clr_hex = f'{color[0]:02X}{color[1]:02X}{color[2]:02X}'
        for old in ser_xml.findall(qn('c:marker')):
            ser_xml.remove(old)
        m = ET.SubElement(ser_xml, qn('c:marker'))
        sym = ET.SubElement(m, qn('c:symbol')); sym.set('val', marker)
        sz_el = ET.SubElement(m, qn('c:size')); sz_el.set('val', '7')
        sp = ET.SubElement(m, qn('c:spPr'))
        fl = ET.SubElement(sp, qn('a:solidFill'))
        rg = ET.SubElement(fl, qn('a:srgbClr')); rg.set('val', clr_hex)
        ln = ET.SubElement(sp, qn('a:ln'))
        lf = ET.SubElement(ln, qn('a:solidFill'))
        lr = ET.SubElement(lf, qn('a:srgbClr')); lr.set('val', clr_hex)
        for sm in ser_xml.findall(qn('c:smooth')):
            ser_xml.remove(sm)
        smooth = ET.SubElement(ser_xml, qn('c:smooth'))
        smooth.set('val', '0')

    # Both axes are value axes in XY scatter; python-pptx returns the
    # X axis through .category_axis (wrapped as a ValueAxis since
    # catAx_lst is empty) and the Y axis through .value_axis.
    x_ax = chart.category_axis
    x_ax.tick_labels.font.name = "Calibri"
    x_ax.tick_labels.font.size = Pt(10)
    x_ax.tick_labels.font.color.rgb = NAVY
    if x_max is not None:
        x_ax.minimum_scale = x_min
        x_ax.maximum_scale = x_max
    if x_unit is not None:
        x_ax.major_unit = x_unit
    x_ax.has_title = True
    x_ax.axis_title.text_frame.text = x_title
    ar = x_ax.axis_title.text_frame.paragraphs[0].runs[0]
    ar.font.name = "Calibri"; ar.font.size = Pt(12)
    ar.font.bold = True; ar.font.italic = True; ar.font.color.rgb = NAVY

    y_ax = chart.value_axis
    y_ax.tick_labels.font.name = "Calibri"
    y_ax.tick_labels.font.size = Pt(10)
    y_ax.tick_labels.font.color.rgb = NAVY
    if y_max is not None:
        y_ax.minimum_scale = y_min
        y_ax.maximum_scale = y_max
    if y_unit is not None:
        y_ax.major_unit = y_unit
    y_ax.has_title = True
    y_ax.axis_title.text_frame.text = y_title
    ar = y_ax.axis_title.text_frame.paragraphs[0].runs[0]
    ar.font.name = "Calibri"; ar.font.size = Pt(12)
    ar.font.bold = True; ar.font.italic = True; ar.font.color.rgb = NAVY

    _add_dashed_gridlines(x_ax._element)
    _add_dashed_gridlines(y_ax._element)

    # Legend
    if legend and legend_pos is not None:
        leg_el = chart.legend._element
        chart.legend.font.name = "Calibri"
        chart.legend.font.size = Pt(11)
        chart.legend.font.color.rgb = NAVY
        chart.legend.include_in_layout = False
        for old in leg_el.findall(qn('c:layout')):
            leg_el.remove(old)
        for old in leg_el.findall(qn('c:legendPos')):
            leg_el.remove(old)
        pos = ET.SubElement(leg_el, qn('c:legendPos')); pos.set('val', 'tr')
        leg_el.remove(pos); leg_el.insert(0, pos)
        layout = ET.Element(qn('c:layout'))
        ml = ET.SubElement(layout, qn('c:manualLayout'))
        ET.SubElement(ml, qn('c:xMode')).set('val', 'edge')
        ET.SubElement(ml, qn('c:yMode')).set('val', 'edge')
        ET.SubElement(ml, qn('c:x')).set('val', legend_pos[0])
        ET.SubElement(ml, qn('c:y')).set('val', legend_pos[1])
        ET.SubElement(ml, qn('c:w')).set('val', legend_pos[2])
        ET.SubElement(ml, qn('c:h')).set('val', legend_pos[3])
        pos.addnext(layout)
        for old in leg_el.findall(qn('c:spPr')):
            leg_el.remove(old)
        leg_spPr = ET.Element(qn('c:spPr'))
        sf = ET.SubElement(leg_spPr, qn('a:solidFill'))
        clr = ET.SubElement(sf, qn('a:srgbClr')); clr.set('val', 'FFFFFF')
        ln = ET.SubElement(leg_spPr, qn('a:ln')); ln.set('w', '6350')
        lf = ET.SubElement(ln, qn('a:solidFill'))
        lc = ET.SubElement(lf, qn('a:srgbClr')); lc.set('val', '0B2B4E')
        layout.addnext(leg_spPr)

    return chart_shape


def _omml_acc_overline(symbol):
    """Inline OMML accent (overline / bar) on a math symbol.

    symbol: e.g. 'K' or 'L' – the variable to wear the bar.
    Returns an OMML fragment intended to be embedded inside an <m:oMath>.
    """
    return (
        '<m:acc>'
          '<m:accPr><m:chr m:val="̅"/></m:accPr>'
          '<m:e>'
            '<m:r>'
              '<a:rPr lang="en-US" b="0" i="1">'
                '<a:latin typeface="Cambria Math"/>'
              '</a:rPr>'
              f'<m:t>{symbol}</m:t>'
            '</m:r>'
          '</m:e>'
        '</m:acc>'
    )


def _add_mixed_textbox(slide, left, top, width, height, segments, *,
                        align=PP_ALIGN.LEFT, default_color=NAVY,
                        default_size=24,
                        margin_left=None, margin_right=None,
                        margin_top=None, margin_bottom=None):
    """Build a textbox whose paragraphs mix plain text runs and inline OMML.

    segments: list of (kind, content, opts) tuples, with kind ∈ {"text",
    "omml", "break"}.  "break" inserts a new paragraph.  Opts may set
    `size`, `bold`, `italic`, `color`, `font` per run.
    """
    box = slide.shapes.add_textbox(int(left), int(top),
                                     int(width), int(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05) if margin_left is None else margin_left
    tf.margin_right = Inches(0.05) if margin_right is None else margin_right
    tf.margin_top = Inches(0.0) if margin_top is None else margin_top
    tf.margin_bottom = Inches(0.0) if margin_bottom is None else margin_bottom

    align_attr = ''
    if align == PP_ALIGN.CENTER: align_attr = ' algn="ctr"'
    elif align == PP_ALIGN.RIGHT: align_attr = ' algn="r"'

    def _start_para():
        return [f'<a:p xmlns:a="{A_NS}" xmlns:m="{M_NS}" xmlns:a14="{A14_NS}">',
                f'<a:pPr{align_attr}/>' if align_attr else '']

    paragraphs = [_start_para()]
    for kind, content, opts in segments:
        if kind == 'break':
            paragraphs[-1].append('<a:endParaRPr lang="en-US"/></a:p>')
            paragraphs.append(_start_para())
            continue
        size_pt = int(opts.get('size', default_size) * 100)
        color = opts.get('color', default_color)
        clr_hex = f'{color[0]:02X}{color[1]:02X}{color[2]:02X}'
        bold_attr = ' b="1"' if opts.get('bold') else ''
        italic_attr = ' i="1"' if opts.get('italic') else ''
        font = opts.get('font', 'Calibri')
        if kind == 'text':
            paragraphs[-1].append(
                f'<a:r><a:rPr lang="en-US" sz="{size_pt}"{bold_attr}{italic_attr}>'
                f'<a:solidFill><a:srgbClr val="{clr_hex}"/></a:solidFill>'
                f'<a:latin typeface="{font}"/>'
                f'</a:rPr><a:t>{content}</a:t></a:r>'
            )
        elif kind == 'omml':
            paragraphs[-1].append(
                f'<a14:m><m:oMath>{content}</m:oMath></a14:m>'
            )
    paragraphs[-1].append('<a:endParaRPr lang="en-US"/></a:p>')
    full_xml = ''.join(''.join(p) for p in paragraphs)

    txBody = tf._txBody
    for old in list(txBody.findall(qn('a:p'))):
        txBody.remove(old)
    # We have to parse each <a:p> separately so the namespaces resolve
    for p_str in full_xml.split('</a:p>'):
        if not p_str.strip(): continue
        p_xml = p_str + '</a:p>'
        new_p = ET.fromstring(p_xml)
        txBody.append(new_p)
        # Apply size+color to any OMML m:r elements inside this paragraph.
        # Color is set to ``default_color`` ONLY when no per-run solidFill
        # is already present — this lets callers tint individual OMML runs
        # via the optional ``color=`` argument on ``_omml_run`` / ``_omml_text``
        # (e.g., the green ΔL / ΔQ in the slide-14 Convention box) without
        # being silently overridden here.
        clr_hex = f'{default_color[0]:02X}{default_color[1]:02X}{default_color[2]:02X}'
        for r in new_p.iter(qn('m:r')):
            arPr = r.find(qn('a:rPr'))
            if arPr is None:
                arPr = ET.Element(qn('a:rPr'))
                arPr.set('lang', 'en-US')
                r.insert(0, arPr)
            arPr.set('sz', str(int(default_size * 100)))
            if arPr.find(qn('a:solidFill')) is None:
                # fill BEFORE a:latin (schema order) or the color is ignored
                sf = arPr.makeelement(qn('a:solidFill'), {})
                srgb = ET.SubElement(sf, qn('a:srgbClr'))
                srgb.set('val', clr_hex)
                arPr.insert(0, sf)
    return box


def _add_math_equation(slide, left, top, width, height, omml_content, *,
                       size_pt=32, color=NAVY, fill=None, line=None,
                       rounded=False, shadow=False, corner_pct=25000):
    """Place an OMML equation in a textbox on the slide.

    omml_content: a string built from _omml_* helpers (without the outer
    <m:oMathPara> wrapper).

    ``rounded=True`` gives the fill box rounded corners (``corner_pct`` in
    OOXML adj units, 25000 ≈ 25%); ``shadow=True`` adds a soft drop
    shadow.  Both default off.  IMPORTANT: the roundRect prstGeom must be
    inserted right after <a:xfrm> and BEFORE the fill, or PowerPoint
    silently drops the shape (see slide-54 fix).
    """
    left, top, width, height = int(left), int(top), int(width), int(height)
    box = slide.shapes.add_textbox(left, top, width, height)

    if fill is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
    if line is not None:
        box.line.color.rgb = line
        box.line.width = Pt(0.75)

    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Replace the default <a:p> with one that hosts a math zone (<a14:m>)
    # containing the oMathPara.  The <a14:m> wrapper is REQUIRED by
    # PowerPoint to recognise OMML inside a textbox – without it PPT just
    # shows empty boxes.
    txBody = tf._txBody
    for p in list(txBody.findall(qn('a:p'))):
        txBody.remove(p)

    sz = int(size_pt * 100)
    clr_hex = '{:02X}{:02X}{:02X}'.format(color[0], color[1], color[2])

    # Use unique namespace prefixes; lxml will resolve them when parsing.
    p_xml = (
        f'<a:p xmlns:a="{A_NS}" xmlns:m="{M_NS}" xmlns:a14="{A14_NS}">'
        f'<a:pPr algn="ctr">'
        f'<a:defRPr sz="{sz}" b="0" i="1">'
        f'<a:solidFill><a:srgbClr val="{clr_hex}"/></a:solidFill>'
        f'<a:latin typeface="Cambria Math"/>'
        f'</a:defRPr>'
        f'</a:pPr>'
        f'<a14:m>'
        f'<m:oMathPara>'
        f'<m:oMathParaPr><m:jc m:val="centerGroup"/></m:oMathParaPr>'
        f'<m:oMath>{omml_content}</m:oMath>'
        f'</m:oMathPara>'
        f'</a14:m>'
        f'<a:endParaRPr lang="en-US" sz="{sz}"/>'
        f'</a:p>'
    )

    new_p = ET.fromstring(p_xml)
    txBody.append(new_p)

    # Set size and color on every OMML run's <a:rPr> so the text renders at
    # the right size.  (The defRPr above is the fallback.)
    for r in new_p.iter(qn('m:r')):
        arPr = r.find(qn('a:rPr'))
        if arPr is None:
            arPr = ET.Element(qn('a:rPr'))
            r.insert(0, arPr)
        arPr.set('sz', str(sz))
        if arPr.get('lang') is None:
            arPr.set('lang', 'en-US')
        # respect per-run colors set via _omml_run/_omml_text(color=...);
        # only runs WITHOUT a fill get the equation's default color
        # (was clobbering firm-color runs — Nico 2026-08-07)
        if not arPr.findall(qn('a:solidFill')):
            sf = ET.Element(qn('a:solidFill'))
            srgb = ET.SubElement(sf, qn('a:srgbClr'))
            srgb.set('val', clr_hex)
            arPr.insert(0, sf)

    # Optional rounded corners + drop shadow on the fill box.  prstGeom
    # MUST sit right after <a:xfrm> and before <a:solidFill>, else PPT
    # silently refuses to render the shape (slide-54 lesson).
    if rounded or shadow:
        spPr = box._element.find(qn('p:spPr'))
        if spPr is not None:
            if rounded:
                for old in spPr.findall(qn('a:prstGeom')):
                    spPr.remove(old)
                prstGeom = ET.Element(qn('a:prstGeom'))
                prstGeom.set('prst', 'roundRect')
                avLst = ET.SubElement(prstGeom, qn('a:avLst'))
                gd = ET.SubElement(avLst, qn('a:gd'))
                gd.set('name', 'adj'); gd.set('fmla', f'val {int(corner_pct)}')
                xfrm = spPr.find(qn('a:xfrm'))
                if xfrm is not None:
                    xfrm.addnext(prstGeom)
                else:
                    spPr.insert(0, prstGeom)
            if shadow:
                _add_drop_shadow(box)
    return box







# ==========================================================================
#  MODULE 1 — BASIC CONCEPTS AND ECONOMIC PRINCIPLES  (everything below
#  is M1; the helper layer above is carried verbatim from Module 2 /
#  Module 7 / Module 3.)
#
#  Pipeline (rerunnable, Module 7 pattern):
#    python _build_Module1.py          -> full 84-slide deck, stubs for the
#                                         3 PollEv slides
#    python _splice_media.py           -> verbatim OOXML splice of the polls
#    python _group_pass.py             -> box+text / shade+frame / pic+caption
#    python _animate.py all apply      -> fade builds per slide plans
# ==========================================================================

import uuid

CREAM = RGBColor(0xFD, 0xF6, 0xE6)
RED = RGBColor(0xC0, 0x00, 0x00)           # source red (C00000)
RED_FF = RGBColor(0xFF, 0x00, 0x00)        # source bright red (FF0000)
GREEN_DK = RGBColor(0x00, 0x7A, 0x33)      # the deck's shift green
                                           # (2026-08-23, Nico: the
                                           # source bright green was
                                           # too light everywhere)
BLUE_PED = RGBColor(0x00, 0x70, 0xC0)      # source blue (S', excess supply)
STEEL = RGBColor(0x95, 0xB3, 0xD7)         # source light-blue supply curve
DARKRED = RGBColor(0xA2, 0x16, 0x2A)       # source MC bar red
NB_BLUE = RGBColor(0x2E, 0x5B, 0x9F)       # net-benefit bar blue

FOOTER_TEXT = "Management 405  ·  Module 1  ·  Basic Concepts and Economic Principles"

TAG_INTRO    = "Module 1 · Introduction"
TAG_LOG      = "Module 1 · Logistics"
TAG_MODELS   = "Module 1 · Economic Models"
TAG_ROADMAP  = "Module 1 · Course Roadmap"
TAG_OUTLINE  = "Module 1 · Outline"
TAG_MARKETS  = "Module 1 · In Class · Markets"
TAG_SD       = "Module 1 · In Class · Supply and Demand"
TAG_OPP      = "Module 1 · In Class · Opportunity Costs"
TAG_SUNK     = "Module 1 · In Class · Sunk Costs"
TAG_CBA      = "Module 1 · In Class · Cost-Benefit and Marginal Analysis"
TAG_WRAP     = "Module 1 · Wrap-Up"
TAG_BACKUP   = "Module 1 · Backup"
TAG_V1       = "Module 1 · Video 1 · Introduction"
TAG_V2       = "Module 1 · Video 2 · Markets"
TAG_V3       = "Module 1 · Video 3 · Supply and Demand"
TAG_V4       = "Module 1 · Video 4 · Market Equilibrium"

STUB_POLL = "PollEverywhere slide — spliced verbatim by _splice_media.py"


# --------------------------------------------------------------------------
# Live slide-number field footer (M7/M2 override — later def wins)
# --------------------------------------------------------------------------

def _add_slidenum_field(slide, left, top, width, height, page_num, *,
                        size=12, color=GRAY):
    tb = slide.shapes.add_textbox(int(left), int(top), int(width), int(height))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    p_el = p._p
    guid = str(uuid.uuid5(uuid.NAMESPACE_DNS, f"m1rev-slidenum-{page_num}")).upper()
    fld = p_el.makeelement(qn('a:fld'), {'id': '{%s}' % guid, 'type': 'slidenum'})
    rPr = p_el.makeelement(qn('a:rPr'),
                           {'lang': 'en-US', 'sz': str(size * 100), 'dirty': '0'})
    fill = p_el.makeelement(qn('a:solidFill'), {})
    clr = p_el.makeelement(qn('a:srgbClr'), {'val': str(color)})
    fill.append(clr)
    rPr.append(fill)
    latin = p_el.makeelement(qn('a:latin'), {'typeface': 'Calibri'})
    rPr.append(latin)
    t = p_el.makeelement(qn('a:t'), {})
    t.text = str(page_num)
    fld.append(rPr)
    fld.append(t)
    p_el.append(fld)
    return tb


def _draw_footer(slide, footer_text, page_num):  # noqa: F811 — field override
    _add_rect(slide, 0, Inches(7.15), SLIDE_W, Inches(0.02), RULE)
    _add_rect(slide, MARGIN, Inches(7.135), GOLD_W, Inches(0.05), GOLD)
    _add_text(slide, MARGIN, Inches(7.20), Inches(11), Inches(0.32),
              footer_text, size=12, color=GRAY)
    _add_slidenum_field(slide, Inches(12.55), Inches(7.20), Inches(0.55),
                        Inches(0.32), page_num)


# --------------------------------------------------------------------------
# Source-image loader (assets extracted by _setup_assets.py)
# --------------------------------------------------------------------------

def _add_media_image(slide, fname, *, left, top, width=None, height=None,
                     rounded=True, shadow=True, corner_pct=8):
    """Place a source-deck image by filename from _source_images/.
    Logos / screenshots / clippings: rounded=False, shadow=False."""
    path = SRC_IMG_DIR / fname
    kwargs = {"left": int(left), "top": int(top)}
    if width is not None:
        kwargs["width"] = int(width)
    if height is not None:
        kwargs["height"] = int(height)
    pic = slide.shapes.add_picture(str(path), **kwargs)
    if rounded:
        _apply_picture_style(pic, corner_pct=corner_pct)
    elif shadow:
        _add_drop_shadow(pic)
    return pic


# --------------------------------------------------------------------------
# Native styled table (navy header, white/cream body) on shadowed backing
# --------------------------------------------------------------------------

def _set_cell_borders(cell, *, color=RULE, weight_pt=0.75):
    tcPr = cell._tc.get_or_add_tcPr()
    for tag in ('a:lnL', 'a:lnR', 'a:lnT', 'a:lnB'):
        for old in tcPr.findall(qn(tag)):
            tcPr.remove(old)
    for tag in reversed(('a:lnL', 'a:lnR', 'a:lnT', 'a:lnB')):
        ln = tcPr.makeelement(qn(tag), {'w': str(int(weight_pt * 12700)),
                                        'cap': 'flat'})
        fill = ln.makeelement(qn('a:solidFill'), {})
        c = fill.makeelement(qn('a:srgbClr'), {'val': str(color)})
        fill.append(c)
        ln.append(fill)
        tcPr.insert(0, ln)


def _add_styled_table(slide, left, top, width, height, rows_data, *,
                      col_widths=None, row_heights=None, font_size=18,
                      header_size=None, backing_pad=Inches(0.15),
                      first_col_bold=True, first_col_align_left=True,
                      cell_fills=None, cell_text_colors=None):
    """rows_data: list of rows (row 0 = navy header). cell_fills /
    cell_text_colors: optional {(r, c): RGBColor} overrides."""
    header_size = header_size or font_size
    left, top, width, height = int(left), int(top), int(width), int(height)
    _add_graphicframe_shadow(slide, left - int(backing_pad),
                             top - int(backing_pad),
                             width + 2 * int(backing_pad),
                             height + 2 * int(backing_pad))
    n_rows, n_cols = len(rows_data), len(rows_data[0])
    gf = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = gf.table
    tblPr = tbl._tbl.find(qn('a:tblPr'))
    if tblPr is not None:
        tblPr.set('firstRow', '0')
        tblPr.set('bandRow', '0')
        for child in list(tblPr):
            tblPr.remove(child)
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = int(w)
    if row_heights:
        for i, h in enumerate(row_heights):
            tbl.rows[i].height = int(h)
    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.margin_left = Inches(0.10)
            cell.margin_right = Inches(0.10)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            hdr = (r == 0)
            if cell_fills and (r, c) in cell_fills:
                cell.fill.solid()
                cell.fill.fore_color.rgb = cell_fills[(r, c)]
            elif hdr:
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else CREAM
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            if first_col_align_left and c == 0 and not hdr:
                p.alignment = PP_ALIGN.LEFT
            else:
                p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(val)
            run.font.name = "Calibri"
            run.font.size = Pt(header_size if hdr else font_size)
            run.font.bold = hdr or (first_col_bold and c == 0)
            if cell_text_colors and (r, c) in cell_text_colors:
                run.font.color.rgb = cell_text_colors[(r, c)]
            else:
                run.font.color.rgb = WHITE if hdr else NAVY
            _set_cell_borders(cell)
    return gf


# --------------------------------------------------------------------------
# Handoff-group injection + Nico's hand-tuned Poll Break badge (from M2)
# --------------------------------------------------------------------------

def _inject_handoff_group(slide, fname, id_base=9500):
    xml = (OUT_DIR / fname).read_text(encoding='utf-8')
    el = ET.fromstring(xml)
    for i, nv in enumerate(el.iter(qn('p:cNvPr'))):
        nv.set('id', str(id_base + i))
    slide.shapes._spTree.append(el)
    return el


def _add_pollbreak_badge(slide):
    """Nico's hand-tuned Poll Break badge (2026-08-15, from Module 2):
    gold parallelogram + navy label, grouped, bottom-right IN FRONT of
    the footer. Call AFTER _draw_footer."""
    _inject_handoff_group(slide, "_handoff_pollbreak.xml", id_base=9600)


# --------------------------------------------------------------------------
# Stubs for spliced poll slides
# --------------------------------------------------------------------------

def make_stub(prs, page_num, section_tag, title, note, *, hidden=False):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, section_tag)
    _draw_action_title(slide, title)
    _add_text(slide, MARGIN, Inches(3.4), RULE_W, Inches(0.8),
              "[ %s ]" % note, size=20, italic=True, color=GRAY,
              font="Calibri", align=PP_ALIGN.CENTER)
    _draw_footer(slide, FOOTER_TEXT, page_num)
    if hidden:
        slide._element.set('show', '0')
    return slide


# --------------------------------------------------------------------------
# SimpleFig — logical→slide transform for the native S/D charts
# --------------------------------------------------------------------------

class SimpleFig:
    def __init__(self, left_in, bottom_in, w_in, h_in, xmax, ymax):
        self.l, self.b, self.w, self.h = left_in, bottom_in, w_in, h_in
        self.xmax, self.ymax = xmax, ymax

    def x(self, xv):
        return Inches(self.l + self.w * xv / self.xmax)

    def y(self, yv):
        return Inches(self.b - self.h * yv / self.ymax)


def _fig_axes(slide, fig, *, weight_pt=2.0,
              x_title="Quantity", y_title="Price ($)", label_size=18):
    """Navy arrow axes + axis titles (y above the top arrow, x right of
    the right arrow, per the source-chart convention)."""
    _add_arrow(slide, (fig.x(0), fig.y(0)),
               (fig.x(0), Inches(fig.b - fig.h - 0.18)),
               color=NAVY, weight_pt=weight_pt, head=True)
    _add_arrow(slide, (fig.x(0), fig.y(0)),
               (Inches(fig.l + fig.w + 0.18), fig.y(0)),
               color=NAVY, weight_pt=weight_pt, head=True)
    if y_title:
        _add_text(slide, Inches(fig.l - 0.75), Inches(fig.b - fig.h - 0.62),
                  Inches(2.0), Inches(0.32), y_title, size=label_size,
                  bold=True, italic=True, color=NAVY, font="Calibri")
    if x_title:
        _add_text(slide, Inches(fig.l + fig.w - 0.4), Inches(fig.b + 0.10),
                  Inches(1.8), Inches(0.32), x_title, size=label_size,
                  bold=True, italic=True, color=NAVY, font="Calibri")


def _fig_line(slide, fig, p0, p1, *, color=NAVY, weight_pt=2.5, dash=None,
              head=False):
    """Straight curve segment in logical coords."""
    return _add_arrow(slide, (fig.x(p0[0]), fig.y(p0[1])),
                      (fig.x(p1[0]), fig.y(p1[1])),
                      color=color, weight_pt=weight_pt, dash=dash, head=head)


def _fig_guide(slide, fig, pt, *, color=GRAY, to_x=True, to_y=True,
               weight_pt=1.25, dash='dash'):
    """Dashed guide lines from a point to both axes. Returns
    (horizontal, vertical) so callers can name/group them."""
    x, y = pt
    h = v = None
    if to_y:
        h = _add_arrow(slide, (fig.x(0), fig.y(y)), (fig.x(x), fig.y(y)),
                       color=color, weight_pt=weight_pt, head=False,
                       dash=dash)
    if to_x:
        v = _add_arrow(slide, (fig.x(x), fig.y(y)), (fig.x(x), fig.y(0)),
                       color=color, weight_pt=weight_pt, head=False,
                       dash=dash)
    return h, v


def _fig_ylab(slide, fig, yv, label, *, color=NAVY, size=18, bold=False,
              width_in=0.95):
    return _add_text(slide, Inches(fig.l - width_in - 0.08),
                     fig.y(yv) - Inches(0.15),
                     Inches(width_in), Inches(0.3), label, size=size,
                     bold=bold, italic=True, color=color, font="Calibri",
                     align=PP_ALIGN.RIGHT)


def _fig_xlab(slide, fig, xv, label, *, color=NAVY, size=18, bold=False):
    return _add_text(slide, fig.x(xv) - Inches(0.5), Inches(fig.b + 0.06),
                     Inches(1.0), Inches(0.3), label, size=size,
                     bold=bold, italic=True, color=color, font="Calibri",
                     align=PP_ALIGN.CENTER)


def _fig_curve_label(slide, fig, xv, yv, label, *, color=NAVY, size=20,
                     bold=True):
    return _add_text(slide, fig.x(xv), fig.y(yv) - Inches(0.16),
                     Inches(0.9), Inches(0.32), label, size=size, bold=bold,
                     italic=True, color=color, font="Calibri")


# --------------------------------------------------------------------------
# Content-slide wrapper: chrome + title + MIDDLE-anchored bullets
# --------------------------------------------------------------------------

def content_slide(prs, page_num, section_tag, title, bullets, *,
                  size=24, sub_size=None, line_spacing_pts=12,
                  bullets_left=None, bullets_width=None,
                  bullets_top=Inches(1.6), bullets_height=Inches(5.35),
                  title_size=None, extras=None, notes=None):
    """M1 standard content slide. Bullet block is MIDDLE-anchored in the
    content region so the text sits vertically centered (course rule)."""
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, section_tag)
    if title_size is None:
        _draw_action_title(slide, title)
    else:
        _add_text(slide, MARGIN, Inches(0.55), RULE_W, Inches(0.7),
                  title, size=title_size, bold=True, color=NAVY,
                  font="Calibri")
        _add_rect(slide, MARGIN, Inches(1.25), RULE_W, Inches(0.02), RULE)
        _add_rect(slide, MARGIN, Inches(1.235), GOLD_W, Inches(0.05), GOLD)
    normalized = [(b, 0) if isinstance(b, str) else b for b in bullets]
    box = _add_hierarchical_bullets(
        slide,
        left=MARGIN if bullets_left is None else bullets_left,
        top=bullets_top,
        width=RULE_W if bullets_width is None else bullets_width,
        height=bullets_height,
        items=normalized,
        size=size, sub_size=sub_size,
        line_spacing_pts=line_spacing_pts,
    )
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    if extras is not None:
        extras(slide)
    _draw_footer(slide, FOOTER_TEXT, page_num)
    if notes:
        _set_notes(slide, notes)
    return slide


# --------------------------------------------------------------------------
# Run post-processing: yellow highlight + hyperlinks on bullet runs
# --------------------------------------------------------------------------

def _highlight_texts(slide, texts, color='FFFF00'):
    """Add <a:highlight> to every run whose text is in *texts*.
    Schema order inside rPr: fill < highlight < latin."""
    remaining = set(texts)
    for shp in slide.shapes:
        if not shp.has_text_frame:
            continue
        for para in shp.text_frame.paragraphs:
            for run in para.runs:
                if run.text in remaining:
                    rPr = run._r.get_or_add_rPr()
                    hl = rPr.makeelement(qn('a:highlight'), {})
                    srgb = hl.makeelement(qn('a:srgbClr'), {'val': color})
                    hl.append(srgb)
                    latin = rPr.find(qn('a:latin'))
                    if latin is not None:
                        latin.addprevious(hl)
                    else:
                        rPr.append(hl)


def _link_runs(slide, mapping):
    """Set external hyperlinks on runs by exact run text.
    mapping: {run_text: url}"""
    for shp in slide.shapes:
        if not shp.has_text_frame:
            continue
        for para in shp.text_frame.paragraphs:
            for run in para.runs:
                if run.text in mapping:
                    run.hyperlink.address = mapping[run.text]


# --------------------------------------------------------------------------
# Problem-Set pointer (generic label per the Module-2 decision)
# --------------------------------------------------------------------------

# --------------------------------------------------------------------------
# Slide-jump link symbol (2026-08-23, Nico)
# --------------------------------------------------------------------------
# The ORIGINAL deck marked its backup jumps with PowerPoint action
# buttons — actionButtonEnd going in, actionButtonBeginning coming back —
# in white fill with a thin black outline. We keep that geometry (the
# glyph is drawn by the preset, so it always points the right way) and
# only recolour it into the deck palette: GOLD fill, NAVY glyph + hairline
# border. The preset draws its glyph in the LINE colour, which is why the
# border is navy rather than gold.
#
# The button IS the click target, so the old 100%-transparent overlay
# rectangles are gone. Those sat on top of the bullet text box and made
# it impossible to select or drag the text in the editor — the symptom
# Nico hit on slide 2.

# The preset shades its glyph from the shape FILL (darkenLess) and only
# strokes it with the line colour, so a WHITE face renders the triangle
# grey, not navy — verified 2026-08-23 against a render. An all-navy
# symbol therefore means a navy face with the glyph knocked out in white.
JUMP_BTN_FILL = NAVY          # backup-jump symbol: navy only (Nico)
JUMP_BTN_GLYPH = WHITE        # glyph knocked out of the navy face
EXT_BTN_FILL = NAVY           # 2026-08-23: external markers go navy too,
EXT_BTN_GLYPH = WHITE         # matching the jump buttons (Nico)


def _add_jump_button(slide, target_slide, *, left, top,
                     width=Inches(0.434), height=Inches(0.210), back=False):
    """Standalone action button carrying a jump-to-slide action.
    back=False -> points right (into the backup section);
    back=True  -> points left (unused since the back pills reverted to
    plain '← Back' text)."""
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ACTION_BUTTON_BEGINNING if back
        else MSO_SHAPE.ACTION_BUTTON_END,
        int(left), int(top), int(width), int(height))
    shp.fill.solid()
    shp.fill.fore_color.rgb = JUMP_BTN_FILL
    shp.line.color.rgb = JUMP_BTN_GLYPH
    shp.line.width = Pt(1.0)
    shp.shadow.inherit = False
    _add_drop_shadow(shp)
    shp.click_action.target_slide = target_slide
    return shp


EXT_LINK_SHAPES = {
    "sound": MSO_SHAPE.ACTION_BUTTON_SOUND,        # podcast / audio
    "document": MSO_SHAPE.ACTION_BUTTON_DOCUMENT,  # article / paper
    "movie": MSO_SHAPE.ACTION_BUTTON_MOVIE,        # video
}


def _add_ext_link_button(slide, kind, *, left, top, url=None,
                         width=Inches(0.60), height=Inches(0.30)):
    """External-link marker in the same action-button family as the
    slide-jump buttons, but keyed to WHAT it opens rather than to a
    direction — so an audio, an article and a video are told apart at a
    glance (2026-08-23, replaces the gold ▶ glyphs)."""
    shp = slide.shapes.add_shape(EXT_LINK_SHAPES[kind],
                                 int(left), int(top),
                                 int(width), int(height))
    shp.fill.solid()
    shp.fill.fore_color.rgb = EXT_BTN_FILL
    shp.line.color.rgb = EXT_BTN_GLYPH
    shp.line.width = Pt(1.25)
    shp.shadow.inherit = False
    _add_drop_shadow(shp)
    if url:
        shp.click_action.hyperlink.address = url
    return shp


def _add_jump_pill(slide, target_slide, *, left, top, width, label,
                   height=Inches(0.5), back=False, fill=None,
                   text_color=None, border=None, size=15):
    """Labelled jump affordance: a rounded pill carrying the label and the
    jump, with the action button seated in a reserved left inset (also
    clickable). Forward pills are white with a gold border and navy text;
    back pills are navy with white text."""
    fill = WHITE if fill is None else fill
    text_color = NAVY if text_color is None else text_color
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  int(left), int(top),
                                  int(width), int(height))
    try:
        pill.adjustments[0] = 0.28
    except Exception:
        pass
    pill.fill.solid()
    pill.fill.fore_color.rgb = fill
    if border is None:
        pill.line.fill.background()
    else:
        pill.line.color.rgb = border
        pill.line.width = Pt(1.5)
    pill.shadow.inherit = False
    _add_drop_shadow(pill)
    btn_w, btn_h = Inches(0.35), Inches(0.196)
    inset = Inches(0.13)
    tf = pill.text_frame
    tf.word_wrap = False
    tf.margin_left = int(inset + btn_w + Inches(0.09))
    tf.margin_right = Inches(0.07)
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = label
    r.font.name = "Calibri"
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = text_color
    pill.click_action.target_slide = target_slide
    btn = _add_jump_button(slide, target_slide,
                           left=int(left + inset),
                           top=int(top + (height - btn_h) / 2),
                           width=btn_w, height=btn_h, back=back)
    return pill, btn


def _add_ps_pointer(slide, *, left, top, label="Problem Set 1",
                    width=Inches(2.5), height=Inches(0.5)):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 int(left), int(top), int(width), int(height))
    try:
        shp.adjustments[0] = 0.28
    except Exception:
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = WHITE
    shp.line.color.rgb = GOLD
    shp.line.width = Pt(1.5)
    shp.shadow.inherit = False
    _add_drop_shadow(shp)
    tf = shp.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
    tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r1 = p.add_run(); r1.text = "➜  "
    r1.font.name = "Calibri"; r1.font.size = Pt(16)
    r1.font.bold = True; r1.font.color.rgb = GOLD
    r2 = p.add_run(); r2.text = label
    r2.font.name = "Calibri"; r2.font.size = Pt(16)
    r2.font.bold = True; r2.font.color.rgb = NAVY
    return shp


# --------------------------------------------------------------------------
# Slide 1 — Title
# --------------------------------------------------------------------------

def slide_01_title(prs):
    slide = _blank_slide(prs)
    # economics comic strip from the original title slide (kept per Nico
    # 2026-08-20; flat clipping with a soft shadow)
    _add_media_image(slide, "ic_s01_rId3.png",
                     left=Inches(2.22), top=162547,
                     # hand-tweaked from Inches(0.45) on 2026-08-23
                     width=Inches(8.90), rounded=False, shadow=True)
    _add_text(slide, Inches(0.9), Inches(3.18), SLIDE_W - Inches(1.8),
              Inches(1.35),
              "Basic Concepts and Economic Principles for Decision Making",
              size=44, bold=True, color=NAVY, font="Calibri",
              align=PP_ALIGN.CENTER)
    _add_text(slide, 0, Inches(4.58), SLIDE_W, Inches(0.6),
              "Module 1",
              size=32, bold=True, color=GOLD, font="Calibri",
              align=PP_ALIGN.CENTER)
    _add_rect(slide, int((SLIDE_W - Inches(4.0)) / 2), Inches(5.36),
              Inches(4.0), 54864, GOLD)
    _add_text(slide, 0, Inches(5.62), SLIDE_W, Inches(0.5),
              "Management 405", size=22, bold=True, color=GRAY,
              font="Calibri", align=PP_ALIGN.CENTER)
    _add_text(slide, 0, Inches(6.16), SLIDE_W, Inches(0.45),
              "Prof. Nico Voigtländer  ·  UCLA Anderson",
              size=18, color=GRAY, font="Calibri", align=PP_ALIGN.CENTER)
    _add_rect(slide, 0, Inches(7.15), SLIDE_W, Inches(0.02), RULE)
    _add_rect(slide, MARGIN, Inches(7.135), GOLD_W, Inches(0.05), GOLD)
    return slide


# --------------------------------------------------------------------------
# Slide 2 — Introduction (was #2)
# --------------------------------------------------------------------------

def slide_02_introduction(prs):
    slide = content_slide(
        prs, 2, TAG_INTRO, "Introduction",
        [
            ("About the instructor", 0, {'bold': True}),
            ("BA in Economics and Environmental Engineering in Berlin/Germany", 1),
            ("MSc in Environmental Engineering and Technology Policy from MIT", 1),
            ("PhD in Economics in Barcelona", 1),
            # the backup jump button is wired in wire_backup_links()
            ("My research: Why are some countries so rich and others so poor?",
             0),
            ("Contact information", 0),
            ([("Email: ", {}),
              ("nico.v@ucla.edu", {'underline': True})], 1, {}),
        ],
    )
    _link_runs(slide, {"nico.v@ucla.edu": "mailto:nico.v@ucla.edu"})
    return slide


# --------------------------------------------------------------------------
# Slides 3–5 — Logistics I–III (was #3–5)
# --------------------------------------------------------------------------

def slide_03_logistics1(prs):
    slide = content_slide(
        prs, 3, TAG_LOG, "Logistics I",
        [
            ("TA for the class:", 0),
            ([("Rafael Rubiao: ", {}),
              ("TA405.EMBA2@gmail.com", {'underline': True})], 1, {}),
            ("PhD student in Economics & Management", 2),
            ("Course website on BruinLearn:", 0),
            ("All course materials posted, including class handouts, slides", 1),
            ("(Optional) Math Review", 1),
            ("Videos", 1),
            ("Classes recorded + livestreamed", 1),
        ],
    )
    _highlight_texts(slide, ["TA405.EMBA2@gmail.com"])
    _link_runs(slide, {"TA405.EMBA2@gmail.com": "mailto:TA405.EMBA2@gmail.com"})
    return slide


def slide_04_logistics2(prs):
    # "Fall 2026" update per Nico 2026-08-20 (was "Fall 2025")
    slide = content_slide(
        prs, 4, TAG_LOG, "Logistics II",
        [
            ("Textbook: Goolsbee, Levitt, Syverson (4th edition)", 0),
            ([("Achieve", {'underline': True}),
              (": Additional practice exercises and video tutorials "
               "(optional, information in the syllabus)", {})], 1, {}),
            ([("Direct link", {'underline': True}),
              (" to Fall 2026 Achieve Site (Course ", {}),
              ("gmechu", {'color': RED_FF}),
              (")", {})], 1, {}),
            ("Weekly readings/podcasts", 0),
            ("Send me suggestions!", 1),
            ("Breaks during class:", 0),
            ("Approx. every 1 – 1.5 hours", 1),
        ],
    )
    _link_runs(slide, {
        "Achieve": "https://achieve.macmillanlearning.com/start",
        "Direct link": "https://achieve.macmillanlearning.com/courses/gmechu",
    })
    return slide


def slide_05_logistics3(prs):
    # Exam-period dates -> [DATE] placeholders per Nico 2026-08-20
    return content_slide(
        prs, 5, TAG_LOG, "Logistics III",
        [
            ([("Overall ", {}), ("5", {'bold': True}),
              (" problem sets", {})], 0, {}),
            ("Due dates on the Class Calendar (on Bruinlearn)", 1),
            ("Study groups pre-assigned", 1),
            ("BruinLearn dropbox or via email to TA", 1),
            ("Questions on problem sets: Email TA", 1),
            ("AI use encouraged for studying!", 1),
            ("Further logistics in Course Syllabus on BruinLearn", 0),
            ("Attendance, participation, grading", 1),
            ("Midterm 35%, Final 40%, Problem Sets 25%", 2),
            ("Midterm Exam Period: [DATE]", 0),
            ("Final Exam Period: [DATE]", 0),
            ("AI tools not allowed during exams", 1,
             {'color': RED_FF}),
        ],
        size=22, sub_size=20, line_spacing_pts=8,
    )


def slide_06_office_hours(prs):
    slide = content_slide(
        prs, 6, TAG_LOG, "Questions and Office Hours",
        [
            ([("TA responsible for", {'color': RED_FF}),
              (" questions related to:", {})], 0, {}),
            ("Problem sets, practice exercises, practice exams", 1),
            ("Achieve", 1),
            ([("Prof. ", {}), ("Nico", {'color': RED_FF}),
              (": All ", {}), ("conceptual", {'color': RED_FF}),
              (" questions (topics from class)", {})], 0, {}),
            ("“Econ & Coffee” chats via Zoom with Nico & TA", 0),
            ("Tuesdays at 4pm + Biweekly on weekends", 1),
            ("Zoom link on BruinLearn", 1),
            ("Talk about concepts from class, then questions about problem sets", 1),
            ("Additional Online TA review sessions each week (recorded)", 0),
            ("TA will send out doodle poll", 1),
        ],
        size=22, sub_size=20, line_spacing_pts=9,
        # hand-tweaked from 1.60/5.35 on 2026-08-20 (Nico moved the box)
        bullets_top=Inches(2.21), bullets_height=Inches(4.13),
    )
    _highlight_texts(slide, ["Tuesdays at 4pm + Biweekly on weekends"])
    return slide


# --------------------------------------------------------------------------
# Slides 7–14 — Why economics, models, philosophy (was #7–14)
# --------------------------------------------------------------------------

def slide_07_why_econ(prs):
    return content_slide(
        prs, 7, TAG_INTRO, "Why Study Economics?",
        [
            ("Make better decisions as managers (and as consumers, "
             "students, etc.)", 0),
            # the backup jump button is wired in wire_backup_links()
            ("To understand, anticipate, contest, and make economic "
             "arguments", 0),
            ("Provides the foundation for most of the MBA program "
             "(e.g., finance, marketing, strategy)", 0),
        ],
        size=28, sub_size=24, line_spacing_pts=18,
        notes=(
            "You should be excited about this course. I think that studying "
            "economics can be useful in a number of respects.\n"
            "I hope that my class will help you make better decisions in your "
            "professional life as well as your personal life.\n"
            "This is a foundational class for the rest of the MBA program. "
            "Many of the other courses that you will take have their "
            "foundations in economics: just to mention a few, strategy, "
            "marketing, human resources, finance and supply chain management. "
            "A solid understanding of economics will help you make the most "
            "of those courses.\n"
            "Last, you will learn how to speak the language of economics. "
            "There are lots of economists in the world, in firms, the "
            "government and the media. It is helpful to understand them, and "
            "to speak their language to communicate with them."),
    )


def slide_08_models(prs):
    return content_slide(
        prs, 8, TAG_MODELS, "Economic Models",
        [
            ([("Models are ", {}),
              ("theoretical frameworks ", {'bold': True}),
              ("that help us understand the world around us.", {})], 0, {}),
            ([("“", {}), ("All models are wrong", {'bold': True}),
              (", but ", {}), ("some are useful", {'bold': True}),
              ("” (George Box – Famous Statistician).", {})], 0, {}),
            ([("Models are ", {}), ("like maps", {'bold': True}),
              (": they ", {}),
              ("simplify the complex world", {'bold': True}),
              (" to help you navigate it.", {})], 0, {}),
            ([("Find a model that ", {}),
              ("matches your needs.", {'bold': True})], 0, {}),
        ],
        size=28, sub_size=24, line_spacing_pts=18,
        notes=(
            "Architects need their rulers. Doctors need their stethoscopes. "
            "Economists need their models.\n"
            "What are economic models? They are theoretical frameworks that "
            "help us understand the world around us. You state some "
            "assumptions, and through a series of logical steps you reach "
            "some conclusions. Sometimes these models are mathematical. "
            "Sometimes they come in the form of charts.\n"
            "Models are always wrong, but sometimes they are useful.\n"
            "In a sense, economic models are like maps. Maps provide a "
            "simplified version of the real world to help us navigate it.\n"
            "You may want to use a simple map, or a more complicated map, "
            "depending on your needs."),
    )


def slide_09_find_model(prs):
    def extras(slide):
        # Three stacked maps (revealed one per question in the build):
        # keep the original paint-over order routes -> city -> neighborhood.
        _add_media_image(slide, "ic_s09_rId3.png",
                         left=Inches(8.30), top=Inches(2.10),
                         width=Inches(4.20), rounded=True)
        _add_media_image(slide, "ic_s09_rId4.png",
                         left=Inches(8.05), top=Inches(2.00),
                         width=Inches(4.60), rounded=True)
        # hand-tweaked 2026-08-23: enlarged from left 8.02" / width 4.65"
        # so the last map fully covers the two painted under it (handout
        # rule — see the Animations section of Teaching/CLAUDE.md)
        _add_media_image(slide, "ic_s09_rId5.png",
                         left=7018559, top=Inches(1.95),
                         width=4566889, rounded=True)

    return content_slide(
        prs, 9, TAG_MODELS, "Find a Model That Matches Your Needs",
        [
            ("Do you need to drive from LA to SF?", 0),
            ("Major routes map", 1),
            ("Do you need to get around LA?", 0),
            ("LA city map", 1),
            ("Do you need to deliver packages in Westwood?", 0),
            ("Neighborhood map", 1),
        ],
        size=26, sub_size=24, line_spacing_pts=20,
        bullets_width=Inches(7.4),
        extras=extras,
        notes=(
            "For example, let’s say you need to drive from Los Angeles to "
            "San Francisco, you probably should use a map that only shows "
            "the major routes. You do not need more detail than that, right? "
            "A more detailed map would be more accurate, but less helpful.\n"
            "If you need to get around Los Angeles, you will probably need a "
            "map that shows every street.\n"
            "And if you need to deliver packages around Los Angeles, you "
            "will need even more detail. For example, you would need to know "
            "the divisions between houses – you do not want to leave the "
            "package in the wrong home."),
    )


def slide_10_homo_economicus(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_INTRO)
    _draw_action_title(slide, "Homo Economicus")

    # Cream panel hosting the two "person" cards
    panel_l, panel_t = Inches(2.35), Inches(1.80)
    panel_w, panel_h = Inches(8.6), Inches(4.35)
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   int(panel_l), int(panel_t),
                                   int(panel_w), int(panel_h))
    try:
        panel.adjustments[0] = 0.05
    except Exception:
        pass
    panel.fill.solid()
    panel.fill.fore_color.rgb = CREAM
    panel.line.color.rgb = NAVY
    panel.line.width = Pt(1.0)
    panel.shadow.inherit = False
    _add_drop_shadow(panel)

    # Row 1: Homo Economicus
    _add_media_image(slide, "ic_s10_rId4.png",
                     left=Inches(2.75), top=Inches(2.05),
                     width=Inches(1.85), rounded=False, shadow=False)
    _add_hierarchical_bullets(
        slide, left=Inches(4.85), top=Inches(2.10),
        width=Inches(5.9), height=Inches(1.9),
        items=[
            ([("Homo Economicus ", {'bold': True}),
              ("(seeks to maximize utility or happiness)", {})], 0,
             {'bullet_style': 'none'}),
            ("100% selfish", 1),
            ("100% rational", 1),
            ("100% informed", 1),
        ],
        size=20, sub_size=18, line_spacing_pts=6)

    # Row 2: Real Human
    _add_media_image(slide, "ic_s10_rId3.png",
                     left=Inches(2.78), top=Inches(4.15),
                     width=Inches(1.80), rounded=False, shadow=False)
    _add_hierarchical_bullets(
        slide, left=Inches(4.85), top=Inches(4.20),
        width=Inches(5.9), height=Inches(1.85),
        items=[
            ("Real Human", 0, {'bold': True, 'bullet_style': 'none'}),
            ("Some are altruistic", 1),
            ("Some make stupid decisions", 1),
            ("Some are poorly informed", 1),
        ],
        size=20, sub_size=18, line_spacing_pts=6)

    # Podcast logo (flat exception) + label with play glyph
    _add_media_image(slide, "ic_s10_rId5.png",
                     left=Inches(11.75), top=Inches(0.62),
                     width=Inches(1.02), rounded=False, shadow=False)
    box = slide.shapes.add_textbox(int(Inches(0.28)), int(Inches(6.35)),
                                   int(Inches(7.89)), int(Inches(0.45)))
    tf = box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r1 = p.add_run(); r1.text = "Podcast: Who is Homo Economicus?"
    r1.font.name = "Calibri"; r1.font.size = Pt(20)
    r1.font.bold = True; r1.font.color.rgb = NAVY
    # 2026-08-23 (Nico): the backup link owns the lower-right corner, so
    # the podcast link is centred in the space left of it. Box centre
    # 4.225" puts the 4.208"-wide label (PIL / Calibri Bold 20 pt) at
    # 2.121"-6.329", and its Sound button at 6.449"-7.049" — clear of the
    # backup pill, which starts at 8.99". No URL in the source deck, so
    # the button is a marker until Nico supplies the podcast link.
    _add_ext_link_button(slide, "sound",
                         left=Inches(6.449), top=Inches(6.421))

    _draw_footer(slide, FOOTER_TEXT, 10)
    _set_notes(slide, (
        "In this course, I will teach you some economic models that are "
        "remarkably wrong, but surprisingly useful. We will abstract away "
        "from a lot of the nuance in the real world, to focus on what "
        "matters the most for the decision at hand. So, how inaccurate can "
        "these economic models be?\n"
        "Well, economic models are based on a mythical creature called homo "
        "economicus. Homo economicus are 100% Selfish: they don’t care how "
        "their actions will affect others. Homo economicus are 100% "
        "Rational: everyone is super smart and can make optimal decisions "
        "in a blink of an eye. Homo economicus are 100% Informed: everyone "
        "knows all the relevant information, including prices, wages, "
        "taxes, and so on.\n"
        "We know real people are not like that. Some people are altruistic. "
        "Some people make stupid decisions – at least I do. Some people are "
        "poorly informed – the world is complicated, so it would be "
        "impossible to stay informed about everything all the time.\n"
        "So, even though we know these models are wrong, they are still "
        "useful. They simplify our thinking, to help us make decisions "
        "quickly. You may wonder: why do we even need the homo economicus "
        "in the first place? Why do we need to simplify the world so much?"))
    return slide


def slide_11_hedgehogs(prs):
    def extras(slide):
        # hand-tweaked 2026-08-23: both shrunk (from 4.15" / 4.40") and
        # stacked so fox and hedgehog sit apart. Two different animals, so
        # the later one is NOT meant to cover the earlier one.
        _add_media_image(slide, "ic_s11_rId3.png",
                         left=Inches(7.95), top=Inches(1.55),
                         width=2565636, rounded=False, shadow=False)
        _add_media_image(slide, "ic_s11_rId4.png",
                         left=7109999, top=3698004,
                         width=2725117, rounded=False, shadow=False)

    return content_slide(
        prs, 11, TAG_INTRO, "Economists as Hedgehogs",
        [
            ("“The fox knows many things, but the hedgehog knows one big "
             "thing” (Greek aphorism).", 0),
            ([("Foxes", {'bold': True}), (" have ", {}),
              ("different strategies", {'bold': True}),
              (" for different problems. They are comfortable with "
               "nuance.", {})], 0, {}),
            ([("Hedgehogs", {'bold': True}),
              (", on the other hand, focus on the big picture. They reduce "
               "every problem to ", {}),
              ("one organizing principle", {'bold': True}),
              (".", {})], 0, {}),
        ],
        size=24, sub_size=22, line_spacing_pts=16,
        bullets_width=Inches(7.3),
        extras=extras,
        notes=(
            "The answer is simple: economists are like hedgehogs. Let me "
            "explain. Someone once said:\n"
            "“The fox knows many things, but the hedgehog knows one big "
            "thing.” If you think these are the lyrics of a new Lady Gaga "
            "song, that’s not it. These words were attributed to the Ancient "
            "Greek poet Archilochus, and popularized by the philosopher "
            "Isaiah Berlin. This allegory suggests that there are two types "
            "of people in the world: the foxes, and the hedgehogs.\n"
            "On the one hand, foxes have different strategies for different "
            "problems. They love nuance, and want to spend all the time in "
            "the world to navigate the nitty-gritty of the problem at hand.\n"
            "Hedgehogs, on the other hand, want to reduce every problem to "
            "one organizing principle. They have one universal tool that "
            "they will attempt to use no matter what the problem is. This "
            "has a clear disadvantage: if you don’t pay attention to nuance, "
            "you will make mistakes. However, being a Hedgehog has its "
            "advantages too: you can make decisions quickly and with little "
            "effort, almost instinctively."),
    )


def slide_12_making_most_1(prs):
    def extras(slide):
        _add_media_image(slide, "ic_s12_rId4.png",
                         left=Inches(1.15), top=Inches(4.45),
                         width=Inches(2.60), rounded=False, shadow=False)
        _add_media_image(slide, "ic_s12_rId5.png",
                         left=Inches(4.00), top=Inches(4.95),
                         width=Inches(1.87), rounded=False, shadow=False)
        _add_media_image(slide, "ic_s12_rId3.png",
                         left=Inches(6.15), top=Inches(4.40),
                         width=Inches(2.84), rounded=False, shadow=False)
        _add_media_image(slide, "ic_s12_rId6.png",
                         left=Inches(9.30), top=Inches(4.85),
                         width=Inches(3.04), rounded=False, shadow=False)
        # (the clipboard image already carries the "B.A. and M.A." label)

    slide = content_slide(
        prs, 12, TAG_INTRO, "Making the Most of the Course",
        [
            ("My main challenge: you have different backgrounds, interests, "
             "skills...", 0),
            ("Some of you are Econ majors, others have not taken Econ since "
             "high school", 1),
            ("Some have stronger quant backgrounds than others", 1),
        ],
        size=26, sub_size=24, line_spacing_pts=14,
        bullets_top=Inches(1.7), bullets_height=Inches(2.4),
        extras=extras,
        notes=(
            "My main challenge when teaching this course is that you have "
            "very different sets of skills, interests and needs. On the one "
            "hand, some of you chose Economics as your college Major, and "
            "may have even done a Masters in Economics or a related field. "
            "On the other hand, some of you have not taken an Economics "
            "class since High School, or, even worse, maybe you took a "
            "class in College but you hated it."),
    )
    return slide


def slide_13_making_most_2(prs):
    return content_slide(
        prs, 13, TAG_INTRO, "Making the Most of the Course",
        [
            ("Choose your own adventure!", 0, {'bold': True}),
            ("Set realistic goals and stick to them", 1),
            ("Focus on the intuition", 0),
            ("To become an economic hedgehog, you don’t need to be a "
             "mathematician", 1),
            ([("Math is fairly simple. Do the ", {}),
              ("Math Test", {'bold': True}), (" and ", {}),
              ("Math Review on Bruinlearn", {'bold': True}),
              ("!", {})], 1, {}),
            ([("I will always explain the ", {}),
              ("economic intuition", {'italic': True, 'underline': True}),
              (" behind the math", {})], 1, {}),
            ("Learn from & with your peers", 0),
            ("Also: Actively use AI tools to deepen your insights. "
             "Challenge AI!", 1),
            ("You will not be allowed to use AI in exams", 1),
        ],
        size=24, sub_size=22, line_spacing_pts=10,
        # 2026-08-23: the "explain it to a six-year-old" line was deleted
        # (it is not an Einstein quote); block re-centred by hand
        bullets_top=2095743, bullets_height=3626634,
    )


def slide_14_teaching_philosophy(prs):
    return content_slide(
        prs, 14, TAG_INTRO, "Teaching Philosophy",
        [
            ("Economics takes a lot of critical thinking:", 0),
            ("I don’t care about memorizing formulas", 1),
            ("Real challenge is to figure out the right tools (and to do "
             "so, you need to build economic intuition)", 1),
            ([("YOU", {'bold': True}),
              (" are the expert in your field:", {})], 0, {}),
            ("I can’t solve the problems you face.", 1),
            ("But I can help you ask the right questions.", 1),
            ("I don’t know the answers, but you may.", 1),
        ],
        size=26, sub_size=24, line_spacing_pts=14,
        notes=(
            "This class is different from some of the other classes you "
            "will take in your MBA. The reason is that, deep down, this is "
            "a class about critical thinking. I do not need you to memorize "
            "a formula. You can look up formulas on Google, using your "
            "phone anytime and anywhere. What I want to teach you is how to "
            "figure out what tool you need to use for the problem at hand, "
            "on your own.\n"
            "In the end, I already told you that I’m a hedgehog. You know "
            "your industry a million times better than I do. You know the "
            "facts, the trends, you are the expert, the fox. I cannot solve "
            "the problems you face. I can, however, teach you how to ask "
            "the right questions. I do not know the answers to those "
            "questions, but you may.\n"
            "Having said that, I’m excited to start this journey with you. "
            "And I hope you are excited too."),
    )


# --------------------------------------------------------------------------
# Slide 15 — Agenda for the Class (course roadmap; "we are here" on Part 1)
# Also reused as slide 60 (Video 1).
# --------------------------------------------------------------------------

def make_roadmap(prs, page_num, *, tag=None):
    """Course roadmap in the Module-3 standard diamond format. Module 1
    (top box) is current: navy box, gold left-arrow 'we are here'."""

    def draw(slide):
        box_h = Inches(0.85)
        narrow_w = Inches(4.6)
        wide_w = Inches(8.6)
        gap = Inches(0.3)
        slide_mid = SLIDE_W // 2

        top_x = slide_mid - wide_w // 2
        top_y = Inches(2.0)
        _add_rounded_filled_box(
            slide, top_x, top_y, wide_w, box_h,
            "1. Basic Principles and Economic Way of Thinking",
            fill=NAVY, text_color=WHITE, size=24, bold=True)

        row2_y = Inches(3.65)
        left_x = slide_mid - gap // 2 - narrow_w
        right_x = slide_mid + gap // 2
        _add_rounded_filled_box(slide, left_x, row2_y, narrow_w, box_h,
                                "2. Value and Demand",
                                fill=FADED, text_color=WHITE, size=26,
                                bold=True)
        _add_rounded_filled_box(slide, right_x, row2_y, narrow_w, box_h,
                                "3. Supply and Cost",
                                fill=FADED, text_color=WHITE, size=26,
                                bold=True)

        bot_x = slide_mid - wide_w // 2
        bot_y = Inches(5.5)
        _add_rounded_filled_box(slide, bot_x, bot_y, wide_w, box_h,
                                "4. Markets, Pricing, and Strategy",
                                fill=FADED, text_color=WHITE, size=24,
                                bold=True)

        top_bottom_y = top_y + box_h
        _add_arrow(slide, (top_x + wide_w // 2, top_bottom_y),
                   (left_x + narrow_w // 2, row2_y),
                   color=FADED, weight_pt=3.0, head=True)
        _add_arrow(slide, (top_x + wide_w // 2, top_bottom_y),
                   (right_x + narrow_w // 2, row2_y),
                   color=FADED, weight_pt=3.0, head=True)
        row2_bottom_y = row2_y + box_h
        _add_arrow(slide, (left_x + narrow_w // 2, row2_bottom_y),
                   (bot_x + wide_w // 2, bot_y),
                   color=FADED, weight_pt=3.0, head=True)
        _add_arrow(slide, (right_x + narrow_w // 2, row2_bottom_y),
                   (bot_x + wide_w // 2, bot_y),
                   color=FADED, weight_pt=3.0, head=True)

        # gold "we are here" arrow pointing at box 1 from the left
        arrow_w = Inches(0.6)
        arrow_h = Inches(0.5)
        arrow_left = top_x - arrow_w - Inches(0.12)
        arrow_top = top_y + (box_h - arrow_h) // 2
        _add_arrow_shape(slide, arrow_left, arrow_top, arrow_w, arrow_h,
                         direction="right", fill=GOLD)
        _add_text(slide, arrow_left - Inches(1.55),
                  top_y + (box_h - Inches(0.32)) // 2,
                  Inches(1.45), Inches(0.32),
                  "we are here", size=16, italic=True, bold=True,
                  color=GOLD, font="Calibri", align=PP_ALIGN.RIGHT)

    return make_diagram_slide(prs, page_num, tag or TAG_ROADMAP,
                              "Agenda for the Class", draw)


# --------------------------------------------------------------------------
# Outline of Module 1 — one maker for all 9 outline slides
# --------------------------------------------------------------------------

RED_MW = RGBColor(0xD5, 0x1A, 0x1A)        # MW's ▼ red


M1_OUTLINE = [
    ("Markets",
     "Video 2: what a market is, and how far it extends"),
    ("Demand and supply",
     "Video 3: how buyers and sellers each respond to price"),
    ("Equilibrium",
     "Video 4: where demand meets supply, and what moves it"),
    ("Economic costs include opportunity costs",
     "In class: the value of the best alternative you gave up"),
    ("Ignore sunk costs",
     "In class: money already spent should not drive the next decision"),
    ("Use cost-benefit and marginal analysis",
     "In class: compare the benefit of the next unit with its cost"),
]


def make_m1_outline(prs, page_num, *, tag=None, title="Outline of Module 1",
                   descriptions=False, highlight_idx=None,
                   highlight_set=None, ps_pointer=False):
    """Module outline in the format converged on for Module 2 (2026-08-23,
    Nico): gold 0.58" circle at x 1.15 with a 25 pt bold navy number, the
    item title 25 pt bold navy and its description 22 pt gray beside it.
    Every item RESERVES the description row, so item positions are
    pixel-identical on every agenda slide; the description text shows only
    for the current topic(s), or for all of them when descriptions=True.
    Section agendas put a cream band with a gold border behind the current
    item (no band on the descriptive overview)."""
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, tag or TAG_OUTLINE)
    _draw_action_title(slide, title)

    hi = set()
    if highlight_idx is not None:
        hi.add(highlight_idx)
    if highlight_set:
        hi.update(highlight_set)
    if descriptions:
        hi = set(range(len(M1_OUTLINE)))

    title_h = Inches(0.42)
    desc_h = Inches(0.38)
    gap = Inches(0.11)
    pitch = title_h + desc_h + gap
    total = pitch * len(M1_OUTLINE) - gap
    top = Inches(1.60)
    bottom = Inches(7.02)
    y = int(top + max(0, (bottom - top - total) // 2))

    last_hi_y = None
    for i, (item, desc) in enumerate(M1_OUTLINE):
        if not descriptions and i in hi:
            last_hi_y = y
            band = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, int(Inches(0.90)),
                int(y - Inches(0.06)), int(Inches(12.15)),
                int(title_h + desc_h + Inches(0.10)))
            try:
                band.adjustments[0] = 0.35
            except Exception:
                pass
            band.fill.solid()
            band.fill.fore_color.rgb = CREAM
            band.line.color.rgb = GOLD
            band.line.width = Pt(1.0)
            band.shadow.inherit = False
            _add_drop_shadow(band)
        circ = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, int(Inches(1.15)), int(y + Inches(0.02)),
            int(Inches(0.58)), int(Inches(0.58)))
        circ.fill.solid()
        circ.fill.fore_color.rgb = GOLD
        circ.line.fill.background()
        circ.shadow.inherit = False
        tf = circ.text_frame
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = str(i + 1)
        run.font.size = Pt(25)
        run.font.bold = True
        run.font.color.rgb = NAVY
        run.font.name = "Calibri"
        rows = [([(item[0].upper() + item[1:],
                   {'bold': True, 'size': 25, 'color': NAVY})], 0,
                 {'bullet_style': 'none', 'space_before_pts': 0})]
        if i in hi:
            rows.append(([(desc, {'size': 22, 'color': GRAY})], 0,
                         {'bullet_style': 'none', 'space_before_pts': 0}))
        _add_hierarchical_bullets(
            slide, Inches(2.05), y, Inches(11.0), title_h + desc_h,
            rows, size=25, line_spacing_pts=0)
        y = int(y + pitch)

    _draw_footer(slide, FOOTER_TEXT, page_num)
    if ps_pointer:
        # pointer box bottom-right over the footer, drawn last so it sits
        # in front; raised when the band reaches that corner
        ptr_y = Inches(6.68)
        if last_hi_y is not None and last_hi_y > Inches(5.6):
            ptr_y = Inches(5.30)
        _add_ps_pointer(slide, left=Inches(10.55), top=ptr_y)
    return slide


def slide_16_outline(prs):
    return make_m1_outline(prs, 16, descriptions=True,
                           ps_pointer=True)


def slide_17_outline_now(prs):
    # the three video topics are done; the in-class trio is next
    return make_m1_outline(prs, 17, highlight_set={3, 4, 5})


# --------------------------------------------------------------------------
# Slides 18–20 — market definition (was #18–20)
# --------------------------------------------------------------------------

def slide_18_recall_market_def(prs):
    return content_slide(
        prs, 18, TAG_MARKETS, "Recall from Video 2: Market Definition",
        [
            ("Extent of a market – A Simple Test", 0,
             {'bold': True, 'bullet_style': 'none'}),
            ("Which products belong to a market?", 0),
            ("Ask: If the price of another product changes, will demand "
             "for your product change?", 0),
            ("Relevant for antitrust litigation (in mergers & "
             "acquisitions)", 0),
        ],
        size=28, sub_size=24, line_spacing_pts=18,
    )


def slide_19_adm(prs):
    return content_slide(
        prs, 19, TAG_MARKETS, "Market Definition Mini-Case: ADM",
        [
            ("In 1990s: Archer-Daniels-Midland Company (ADM) wanted to "
             "acquire Clinton Corn Processing Company (CCP).", 0),
            ([("Both had substantial market shares in ", {}),
              ("corn syrup", {'color': RED})], 0, {}),
            ("Department of Justice (DOJ) challenged the acquisition:", 0),
            ("“Would lead to a dominant producer of corn syrup with the "
             "power to push prices above competitive levels”", 1),
            ("ADM fought the DOJ decision in court", 0),
            ([("Basic issue", {'bold': True}),
              (": Is corn syrup a distinct market?", {})], 0, {}),
            ([("ADM’s argument: The market is broader: ", {}),
              ("sweeteners", {'color': RED})], 0, {}),
            ("Includes sugar, corn syrup, other artificial sweeteners", 1),
            ("Used price-test to support their case", 1),
            ("ADM won in court", 1),
        ],
        size=22, sub_size=20, line_spacing_pts=9,
    )


def slide_20_netflix(prs):
    def extras(slide):
        _add_media_image(slide, "ic_s20_rId2.png",
                         left=Inches(7.85), top=Inches(2.45),
                         width=Inches(4.85), rounded=True)
        _add_discussion_break(slide, text="Class Discussion",
                              width=Inches(4.4))

    return content_slide(
        prs, 20, TAG_MARKETS, "Market Definition: the Case of Netflix",
        [
            ("Define Netflix’s market", 0,
             {'bold': True, 'bullet_style': 'none'}),
            ("Online streaming?", 0),
            ("All shows/movies, incl. theatres?", 0),
            ("All entertainment?", 0),
        ],
        size=28, sub_size=24, line_spacing_pts=18,
        bullets_width=Inches(7.0), bullets_height=Inches(4.2),
        extras=extras,
    )


# --------------------------------------------------------------------------
# Slides 21–25 — heatwaves poll + Swiftonomics (MW-5) (was #21–24 + MW #52/55)
# --------------------------------------------------------------------------

def slide_21_heatwaves(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_SD)
    _draw_action_title(slide,
                       "How Can Heatwaves Affect the Demand for ACs?")
    # LA Times headline clipping (flat) + AC-households map
    # both hand-tweaked 2026-08-23 (from top 1.62" / 3.85") so the lower
    # figure stops clear of the footer rule at 7.15"
    _add_media_image(slide, "ic_s21_rId3.png",
                     left=Inches(1.55), top=1288504,
                     width=Inches(10.25), rounded=False, shadow=True)
    _add_media_image(slide, "ic_s21_rId2.png",
                     left=Inches(4.35), top=3202270,
                     width=Inches(4.65), rounded=False, shadow=True)
    _draw_footer(slide, FOOTER_TEXT, 21)
    _add_pollbreak_badge(slide)
    return slide


def slide_22_poll_ac(prs):
    return make_stub(prs, 22, TAG_SD, "Poll: heatwaves and AC demand",
                     STUB_POLL)


def slide_23_swiftonomics(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_SD)
    _draw_action_title(slide, "Swiftonomics")
    # WSJ logo (flat) + the two engagement photos (rounded)
    _add_media_image(slide, "mw_s52_rId3.png",
                     left=Inches(0.95), top=Inches(2.85),
                     width=Inches(3.35), rounded=False, shadow=False)
    _add_media_image(slide, "mw_s52_rId4.jpg",
                     left=Inches(4.85), top=Inches(2.05),
                     width=Inches(3.65), rounded=True)
    _add_media_image(slide, "mw_s52_rId5.png",
                     left=Inches(8.90), top=Inches(1.85),
                     width=Inches(3.75), rounded=True)
    box = _add_hierarchical_bullets(
        slide, left=MARGIN + Inches(0.3), top=Inches(4.95),
        width=Inches(9.5), height=Inches(1.0),
        items=[("How did Taylor Swift’s engagement affect the demand for "
                "diamonds?", 0)],
        size=26, line_spacing_pts=0)
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    _draw_footer(slide, FOOTER_TEXT, 24)
    _add_pollbreak_badge(slide)
    _set_notes(slide, (
        "Adopted from Melanie Wasserman's Module 1 deck (2026): Taylor "
        "Swift announced her engagement to Travis Kelce in August 2025, "
        "posting a picture of her diamond ring. Ask the class what this "
        "does to the demand curve for diamonds, then run the poll. NOTE: "
        "the PollEv activity still carries the old wording about the 2023 "
        "decline in engagements — reword it in PollEverywhere."))
    return slide


def slide_24_poll_diamonds(prs):
    return make_stub(prs, 25, TAG_SD, "Poll: diamond demand", STUB_POLL)


def slide_25_swift_solution(prs):
    def extras(slide):
        _add_media_image(slide, "mw_s55_rId3.png",
                         left=Inches(9.35), top=Inches(2.30),
                         width=Inches(3.10), rounded=True)
        _add_media_image(slide, "mw_s55_rId4.png",
                         left=Inches(9.10), top=Inches(5.35),
                         width=Inches(3.35), rounded=False, shadow=False)

    return content_slide(
        prs, 26, TAG_SD, "Solution: Swiftonomics",
        [
            ([("Taylor Swift’s engagement caused the demand curve to ", {}),
              ("shift to the right", {'bold': True})], 0, {}),
            ("“Google searches for “diamond ring” surged after the "
             "announcement.”", 0),
            ("“Diamond-jewelry retailers […] saw their stock prices rise "
             "[…] on Tuesday, the day Swift announced her engagement with "
             "football star Travis Kelce, posting a picture of her diamond "
             "ring.”", 0),
        ],
        size=24, sub_size=22, line_spacing_pts=18,
        bullets_width=Inches(8.3),
        extras=extras,
        notes=(
            "The engagement is a non-price factor that raises demand at "
            "every price, so the demand curve shifts to the right. The two "
            "quotes are from the Wall Street Journal's coverage the day "
            "after the August 2025 announcement."),
    )


# --------------------------------------------------------------------------
# Slides 26–27 — COVID and the market for tea (was #25–26)
# --------------------------------------------------------------------------

def slide_26_tea(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_SD)
    _draw_action_title(slide,
                       "The Effects of COVID on Supply and Demand for Tea",
                       gold_len=GOLD_W)
    _add_media_image(slide, "ic_s25_rId3.png",
                     left=Inches(0.45), top=Inches(1.60),
                     width=Inches(6.55), rounded=False, shadow=True)
    _add_media_image(slide, "ic_s25_rId4.png",
                     left=Inches(0.45), top=Inches(3.95),
                     width=Inches(4.45), rounded=True)
    _add_media_image(slide, "ic_s25_rId5.png",
                     left=Inches(9.45), top=Inches(1.75),
                     width=Inches(3.25), rounded=False, shadow=True)
    _draw_footer(slide, FOOTER_TEXT, 27)
    _set_notes(slide, (
        "Note: Use this article starting in 2023: "
        "https://www.wsj.com/articles/lithium-prices-soar-turbocharged-by-"
        "electric-vehicle-demand-and-scant-supply-11639334956"))
    return slide


def _sd_chart(slide, fig, curves, points, xlabels, ylabels, arrows=(),
              guide_color=GREEN_DK):
    """Generic native S/D chart: curves = list of (p0, p1, color, dash,
    label, label_pos); points = list of logical intersections that get
    dashed guides; xlabels/ylabels = [(val, text)] axis labels."""
    # 2026-08-23: every shape gets a stable name (sdcurve/sdlabel/
    # sdguide/sdxlab/sdylab/sdarrow) so the grouping pass and the
    # per-slide animation plans can address it without relying on
    # connector indices, which shift as soon as anything is grouped.
    def _nm(shape, name):
        if shape is not None:
            try:
                shape.name = name
            except Exception:
                pass
        return shape

    _fig_axes(slide, fig)
    for i, pt in enumerate(points):
        h, v = _fig_guide(slide, fig, pt, color=guide_color)
        _nm(h, "sdguide:h:%d" % i)
        _nm(v, "sdguide:v:%d" % i)
    for p0, p1, color, dash, label, lpos in curves:
        _nm(_fig_line(slide, fig, p0, p1, color=color, weight_pt=2.75,
                      dash=dash), "sdcurve:%s" % (label or "?"))
        if label:
            _nm(_fig_curve_label(slide, fig, lpos[0], lpos[1], label,
                                 color=color), "sdlabel:%s" % label)
    for val, text in xlabels:
        _nm(_fig_xlab(slide, fig, val, text, size=16), "sdxlab:%s" % text)
    for val, text in ylabels:
        _nm(_fig_ylab(slide, fig, val, text, size=16), "sdylab:%s" % text)
    for i, (a0, a1, color) in enumerate(arrows):
        _nm(_add_arrow(slide, (fig.x(a0[0]), fig.y(a0[1])),
                       (fig.x(a1[0]), fig.y(a1[1])),
                       color=color, weight_pt=2.0, head=True),
            "sdarrow:%d" % i)


def slide_27_tea_market(prs):
    def extras(slide):
        fig = SimpleFig(8.0, 6.35, 4.35, 4.15, 10, 10)
        _sd_chart(
            slide, fig,
            curves=[
                ((1, 8), (8, 1), GOLD, None, "D", (8.05, 1.35)),
                ((2, 8.5), (8.6, 1.9), GOLD, 'dash', "D1", (8.65, 2.2)),
                ((1.5, 1.5), (8.5, 8.5), STEEL, None, "S", (8.55, 8.75)),
                ((0.5, 4), (6, 9.5), STEEL, 'dash', "S1", (5.6, 10.0)),
            ],
            points=[(4.5, 4.5), (3.5, 7.0)],
            xlabels=[(4.75, "Q0"), (3.1, "Q1")],
            ylabels=[(4.5, "P0"), (7.0, "P1")],
            arrows=[((5.8, 5.8), (4.6, 6.9), GRAY),
                    ((6.9, 5.4), (7.6, 4.7), GRAY)],
        )
        # D1: y = 10.5 - x ; S1: y = x + 3.5 ; both meet at (3.5, 7)

    return content_slide(
        prs, 28, TAG_SD, "Market for Tea",
        [
            ([("Demand curve: ", {'bold': True}),
              ("shifted due to increased consumer demand during the "
               "pandemic", {})], 0, {}),
            ([("Supply curve", {'bold': True}),
              (": shifted due to bad weather, labor shortages, port "
               "closures during Covid", {})], 0, {}),
        ],
        size=24, sub_size=22, line_spacing_pts=18,
        bullets_width=Inches(6.2),
        extras=extras,
    )


# --------------------------------------------------------------------------
# Slides 28–31 — disasters + avocados (was #27–30)
# --------------------------------------------------------------------------

def slide_28_disasters(prs):
    def extras(slide):
        _add_media_image(slide, "ic_s27_rId2.jpeg",
                         left=Inches(0.65), top=Inches(2.05),
                         width=Inches(5.85), rounded=True)
        _add_discussion_break(slide, text="Class Discussion",
                              width=Inches(4.4))

    return content_slide(
        prs, 29, TAG_SD, "Shortages when Disasters Loom",
        [
            ("Why are there frequently shortages when major disasters loom "
             "(storms, the Covid pandemic…)?", 0),
            ("What role does the perception of ‘fairness’ play in this "
             "context?", 0),
        ],
        size=26, sub_size=24, line_spacing_pts=20,
        bullets_left=Inches(6.95), bullets_width=Inches(6.0),
        bullets_height=Inches(4.4),
        extras=extras,
    )


def slide_29_avocado_clip(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_SD)
    _draw_action_title(slide, "Mini-Case on Avocados")
    _add_media_image(slide, "ic_s28_rId2.png",
                     left=Inches(2.45), top=Inches(1.70),
                     width=Inches(8.45), rounded=False, shadow=True)
    _draw_footer(slide, FOOTER_TEXT, 30)
    return slide


def slide_30_avocado_bullets(prs):
    def extras(slide):
        _add_media_image(slide, "ic_s29_rId4.png",
                         left=Inches(1.20), top=Inches(4.60),
                         width=Inches(3.30), rounded=True)
        _add_media_image(slide, "ic_s29_rId2.png",
                         left=Inches(4.95), top=Inches(4.45),
                         width=Inches(3.75), rounded=True)
        _add_media_image(slide, "ic_s29_rId3.png",
                         left=Inches(9.15), top=Inches(4.30),
                         width=Inches(2.15), rounded=True)

    return content_slide(
        prs, 31, TAG_SD, "Mini-Case on Avocados",
        [
            ([("From December 2016 to June 2017, price of Mexico’s Hass "
               "avocados ", {}),
              ("more than doubled", {'color': RED})], 0, {}),
            ([("From July 2017 to November 2017, ", {}),
              ("price halved", {'color': RED})], 0, {}),
            ("Can be explained using supply and demand framework", 0),
        ],
        size=26, sub_size=24, line_spacing_pts=12,
        bullets_top=Inches(1.6), bullets_height=Inches(2.7),
        extras=extras,
    )


def slide_31_avocado_market(prs):
    def extras(slide):
        fig = SimpleFig(8.0, 6.35, 4.35, 4.15, 10, 10)
        _sd_chart(
            slide, fig,
            curves=[
                ((1, 8), (8, 1), GOLD, None, "D", (8.05, 1.35)),
                ((4, 9), (9.6, 3.4), GOLD, 'dash', "D1", (9.3, 3.9)),
                ((1.5, 1.5), (8.5, 8.5), STEEL, None, "S", (8.55, 8.75)),
                ((0.5, 4), (5.8, 9.3), STEEL, 'dash', "S1", (5.35, 9.8)),
                ((5, 1), (9.5, 5.5), STEEL, 'dash', "S2", (9.55, 5.85)),
            ],
            points=[(4.5, 4.5), (4.75, 8.25), (8.5, 4.5)],
            xlabels=[(4.1, "Q0"), (5.15, "Q1"), (8.5, "Q2")],
            ylabels=[(8.25, "PPeak"), (4.5, "P2 = P0")],
            arrows=[((3.4, 7.1), (2.5, 8.0), GRAY),
                    ((7.0, 6.6), (7.9, 5.7), GRAY)],
        )
        # D1: y = 13 - x ; S1: y = x + 3.5 (peak at 4.75, 8.25);
        # S2: y = x - 4  (back down: crosses D1 at Q2 = 8.5, P2 = P0)

    return content_slide(
        prs, 32, TAG_SD, "Market for Avocados",
        [
            ("Demand curve: shifts due to consumer craze for guacamole & "
             "avocado toast (D1)", 0),
            ("Supply curve (December 2016-June 2017): shifts due to dry "
             "weather leading to meager crop (S1)", 0),
            ("Supply curve (July 2017-November 2017): shifts back due to "
             "superior pest control and trees entering the high-yielding "
             "half of their two year production cycle (S2)", 0),
        ],
        size=20, sub_size=18, line_spacing_pts=12,
        bullets_width=Inches(6.2),
        extras=extras,
    )


# --------------------------------------------------------------------------
# Slides 32–33 — shocks framework + Ukraine wheat (was #31–32)
# --------------------------------------------------------------------------

def slide_32_steps(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_SD)
    _draw_action_title(slide,
                       "Steps to Analyze “Shocks” to Supply and Demand")
    box = _add_bulleted_list(
        slide, left=MARGIN + Inches(0.4), top=Inches(1.85),
        width=RULE_W - Inches(0.8), height=Inches(4.9),
        items=[
            "Figure out what the shock is",
            "Determine whether the shock shifts the demand or supply curve",
            "Draw the market’s demand and supply curves before and after "
            "the shock",
        ],
        size=28, color=NAVY, bullet_color=NAVY,
        line_spacing_pts=26,
        autonum_scheme='arabicPeriod',
    )
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    _draw_footer(slide, FOOTER_TEXT, 33)
    return slide


def slide_33_wheat(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_SD)
    _add_text(slide, MARGIN, Inches(0.55), RULE_W, Inches(0.7),
              "Example: How Did the War in Ukraine Affect the Price of "
              "Wheat?", size=27, bold=True, color=NAVY, font="Calibri")
    _add_rect(slide, MARGIN, Inches(1.25), RULE_W, Inches(0.02), RULE)
    _add_rect(slide, MARGIN, Inches(1.235), GOLD_W, Inches(0.05), GOLD)

    # FRED wheat-price chart (kept as clipping)
    img_l, img_t, img_w = Inches(0.75), Inches(1.70), Inches(11.85)
    img_h = Inches(11.85 / 9.74 * 3.45)          # aspect from source
    _add_media_image(slide, "ic_s32_rId2.png",
                     left=img_l, top=img_t, width=img_w,
                     rounded=False, shadow=True)

    # Native annotations (positions in image fractions, tuned on render):
    # war-begins line ~ Feb 2022, grain agreement ~ Jul 2022.
    x_war = img_l + int(img_w) * 0.655
    x_grain = img_l + int(img_w) * 0.720
    y_top = img_t + int(img_h) * 0.10
    y_bot = img_t + int(img_h) * 0.92
    _add_arrow(slide, (x_war, y_top), (x_war, y_bot),
               color=RED_FF, weight_pt=2.25, head=False)
    _add_arrow(slide, (x_grain, y_top + int(img_h) * 0.06), (x_grain, y_bot),
               color=GREEN_DK, weight_pt=2.25, head=False)
    _add_text(slide, x_war - Inches(1.85), y_top - Inches(0.05),
              Inches(1.75), Inches(0.3), "War begins", size=14, bold=True,
              color=RED_FF, font="Calibri", align=PP_ALIGN.RIGHT)
    _add_text(slide, x_grain + Inches(0.10), y_top + int(img_h) * 0.06,
              Inches(2.6), Inches(0.55), "Black sea grain agreement",
              size=14, bold=True, color=GREEN_DK, font="Calibri")

    # Callouts under the chart
    box1 = slide.shapes.add_textbox(int(Inches(3.4)), int(Inches(6.15)),
                                    int(Inches(3.4)), int(Inches(0.75)))
    tf = box1.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "Supply shifts left  ➜  P rises"
    r.font.name = "Calibri"; r.font.size = Pt(16); r.font.bold = True
    r.font.color.rgb = RED_FF
    box2 = slide.shapes.add_textbox(int(Inches(8.4)), int(Inches(6.15)),
                                    int(Inches(3.6)), int(Inches(0.75)))
    tf = box2.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "Supply shifts right  ➜  P falls"
    r.font.name = "Calibri"; r.font.size = Pt(16); r.font.bold = True
    r.font.color.rgb = GREEN_DK
    _add_arrow(slide, (Inches(5.1), Inches(6.15)), (x_war, y_bot),
               color=RED_FF, weight_pt=1.5, head=True)
    _add_arrow(slide, (Inches(10.2), Inches(6.15)),
               (x_grain + Inches(0.3), y_bot),
               color=GREEN_DK, weight_pt=1.5, head=True)

    _draw_footer(slide, FOOTER_TEXT, 34)
    return slide


# --------------------------------------------------------------------------
# Slides 34–35 — LA residential real estate (NEW – MW #67/68)
# --------------------------------------------------------------------------

def slide_34_la_case(prs):
    def extras(slide):
        _add_media_image(slide, "mw_s67_rId3.png",
                         left=Inches(1.05), top=Inches(1.95),
                         width=Inches(3.30), rounded=True)
        _add_discussion_break(slide, text="Class Discussion",
                              width=Inches(4.4))

    slide = content_slide(
        prs, 35, TAG_SD, "Mini-Case: Los Angeles Residential Real Estate",
        [
            ("From Sep-2021 to Sep-2025, in Los Angeles (Redfin data):", 0),
            ([("Inflation-adjusted home prices: ", {}),
              ("▼5%", {'bold': True, 'color': RED_MW})], 1, {}),
            ([("Home sales: ", {}),
              ("▼40%", {'bold': True, 'color': RED_MW})], 1, {}),
            ("Explain using supply and demand framework!", 0),
        ],
        size=24, sub_size=22, line_spacing_pts=16,
        bullets_left=Inches(4.85), bullets_width=Inches(7.9),
        bullets_height=Inches(4.2),
        extras=extras,
        notes="Source: https://www.redfin.com/city/11203/CA/Los-Angeles/"
              "housing-market",
    )
    return slide


def slide_35_la_market(prs):
    def extras(slide):
        fig = SimpleFig(4.4, 5.30, 5.0, 3.45, 10, 10)
        _sd_chart(
            slide, fig,
            curves=[
                ((1, 8), (8, 1), GOLD, None, "D0", (8.05, 1.35)),
                ((0.9, 6), (6.4, 0.5), GOLD, 'dash', "D1", (6.15, 1.05)),
                ((1.5, 1.5), (8.5, 8.5), STEEL, None, "S0", (8.55, 8.75)),
                ((0.5, 2), (7.5, 9), STEEL, 'dash', "S1", (7.1, 9.4)),
            ],
            points=[(4.5, 4.5), (2.7, 4.2)],
            xlabels=[(2.7, "Q1"), (4.5, "Q0")],
            ylabels=[(4.9, "P0"), (3.8, "P1")],
            arrows=[((4.0, 6.2), (3.2, 5.4), GOLD),
                    ((6.6, 5.6), (5.8, 6.4), GOLD)],
        )
        # D1: y = 6.9 - x ; S1: y = x + 1.5 ; both meet at (2.7, 4.2)

    slide = content_slide(
        prs, 36, TAG_SD, "Market for Los Angeles Residential Real Estate",
        [
            ([("Demand curve: ", {'bold': True}),
              ("shift inward due to higher mortgage rates and "
               "out-migration", {})], 0, {}),
            ([("Supply curve: ", {'bold': True}),
              ("owners ”locked in” to lower mortgage rates and postponed "
               "selling", {})], 0, {}),
        ],
        size=20, sub_size=18, line_spacing_pts=8,
        bullets_top=Inches(5.75), bullets_height=Inches(1.3),
        bullets_left=MARGIN + Inches(0.5),
        bullets_width=RULE_W - Inches(1.0),
        extras=extras,
        notes=(
            "Since both prices (in real terms) and sales dropped, this "
            "suggests both demand and supply curves shifted left. For "
            "example, perhaps demand fell due to higher mortgage rates and "
            "migration. Maybe supply fell as owners were “locked in” low "
            "mortgages and delayed selling."),
    )
    return slide



# ==========================================================================
#  BATCH C — slides 36–58: opportunity costs, sunk costs, CBA, wrap-up
# ==========================================================================

GREEN_CELL = RGBColor(0x92, 0xD0, 0x50)     # fruit-table value fill
RED_CELL = RGBColor(0xFF, 0x50, 0x50)       # fruit-table opp-cost fill


def slide_36_outline_opp(prs):
    return make_m1_outline(prs, 39, highlight_idx=3, tag=TAG_OPP)


def slide_37_opp_costs(prs):
    def extras(slide):
        pic = _add_media_image(slide, "ic_s34_rId3.png",
                               left=Inches(9.55), top=Inches(0.60),
                               width=Inches(3.10), rounded=False,
                               shadow=False)
        pic.click_action.hyperlink.address = (
            "https://podcasts.apple.com/us/podcast/core-principle-2-the-"
            "opportunity-cost-principle")
        _add_convention_box(
            slide, Inches(2.35), Inches(5.85), Inches(8.6), Inches(0.78),
            runs=[
                ("Full economic cost ", {'bold': True, 'color': RED,
                                         'size': 20}),
                ("= implicit cost + explicit cost", {'color': NAVY,
                                                     'size': 20}),
            ],
            align=PP_ALIGN.CENTER)

    return content_slide(
        prs, 40, TAG_OPP, "Opportunity Costs",
        [
            ([("The ", {}),
              ("opportunity cost", {'underline': True}),
              (" is the ", {}),
              ("value", {'bold': True, 'underline': True}),
              (" of the next best alternative sacrificed", {})], 0, {}),
            ("what you give up to get something", 1),
            ([("The opportunity cost is not always monetary. It is "
               "typically ", {}),
              ("implicit", {'color': RED}),
              (" rather than explicit", {})], 0, {}),
            ("Explicit cost: actual payment", 1),
            ("Implicit cost: value of forgone alternative action", 1),
        ],
        size=26, sub_size=24, line_spacing_pts=14,
        bullets_top=Inches(1.6), bullets_height=Inches(4.0),
        extras=extras,
    )


def slide_38_fruit_table(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_OPP)
    _draw_action_title(slide, "Opportunity Cost: A Simple Example")

    lead = slide.shapes.add_textbox(int(MARGIN), int(Inches(1.62)),
                                    int(RULE_W), int(Inches(0.5)))
    tf = lead.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r1 = p.add_run(); r1.text = "Suppose that I can choose "
    r2 = p.add_run(); r2.text = "one"
    r3 = p.add_run(); r3.text = " of the following:"
    for r in (r1, r2, r3):
        r.font.name = "Calibri"; r.font.size = Pt(24)
        r.font.color.rgb = NAVY
    r2.font.bold = True; r2.font.underline = True

    rows = [
        ("Alternatives", "My Value", "My Opp. Cost"),
        ("Banana", "9", "5"),
        ("Apple", "5", "9"),
        ("Cucumber", "3", "9"),
        ("Dragon fruit", "1", "9"),
    ]
    fills = {}
    text_colors = {}
    for r in range(1, 5):
        fills[(r, 1)] = GREEN_CELL
        fills[(r, 2)] = RED_CELL
        text_colors[(r, 1)] = NAVY
        text_colors[(r, 2)] = NAVY
    _add_styled_table(
        slide, Inches(3.75), Inches(2.40), Inches(5.85), Inches(2.55), rows,
        col_widths=[Inches(2.45), Inches(1.6), Inches(1.8)],
        font_size=20, header_size=20,
        cell_fills=fills, cell_text_colors=text_colors)

    box = _add_hierarchical_bullets(
        slide, left=MARGIN + Inches(0.8), top=Inches(5.45),
        width=RULE_W - Inches(1.6), height=Inches(1.45),
        items=[
            ([("The opportunity cost ", {'color': RED}),
              ("is the ", {}),
              ("value", {'bold': True, 'color': RED_FF}),
              (" of the next best alternative foregone (the lost ", {}),
              ("opportunity", {'italic': True}),
              (")", {})], 0, {}),
            ([("Opportunity cost depends on ", {}),
              ("preferences", {'color': RED}),
              (" – it differs across individuals", {})], 0, {}),
        ],
        size=22, line_spacing_pts=8)
    box.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    _draw_footer(slide, FOOTER_TEXT, 41)
    return slide


def slide_39_present(prs):
    return content_slide(
        prs, 42, TAG_OPP,
        "Example: Opportunity Costs when Buying a Present",
        [
            ("3 possible presents, price of each is $30", 0),
            ("But different values to recipient", 0),
        ],
        size=28, sub_size=24, line_spacing_pts=20,
    )


def slide_40_mba_cost(prs):
    def extras(slide):
        _add_discussion_break(slide, text="Class Discussion",
                              width=Inches(4.4))

    return content_slide(
        prs, 43, TAG_OPP,
        "What is the Full Economic Cost of an MBA Degree?",
        [
            ("Distinguish explicit vs. implicit", 0),
            ("  ", 0), ("  ", 0), ("  ", 0), ("  ", 0), ("  ", 0),
            ("Do remote options change the (opportunity) cost?", 0),
        ],
        size=26, sub_size=24, line_spacing_pts=16,
        title_size=30,
        extras=extras,
    )


def slide_41_flip_house(prs):
    def extras(slide):
        _add_ps_pointer(slide, left=MARGIN + Inches(0.2), top=Inches(6.35))

    slide = content_slide(
        prs, 44, TAG_OPP, "Flip a House or Work?",
        [
            ("Flip house:", 0),
            ("1 year of full-time work", 1),
            ("Buy for 500k", 1),
            ("Invest another 100k (+ your full time)", 1),
            ("Expected sale price in a year: 700k", 1),
            ("Consulting job offer:", 0),
            ("Full-time job (not compatible with flipping the house)", 1),
            ("Pays 150K", 1),
            ([("What is the ", {}),
              ("economic", {'color': RED}),
              (" profit of flipping the house?", {})], 0, {}),
        ],
        size=24, sub_size=22, line_spacing_pts=10,
        bullets_height=Inches(4.7), bullets_top=Inches(1.5),
        extras=extras,
    )
    _add_pollbreak_badge(slide)
    return slide


def slide_42_poll_flip(prs):
    return make_stub(prs, 45, TAG_OPP, "Poll: flip-a-house profit",
                     STUB_POLL)


def slide_43_flip_solution(prs):
    return content_slide(
        prs, 46, TAG_OPP, "Solution: Flip a House or Work?",
        [
            ("Flip the house:", 0),
            ("Economic revenues: 700k", 1),
            ("Explicit costs: 600k", 1),
            ("500k for the purchase, 100k for remodeling", 2),
            ("Implicit (opportunity) costs: 150k", 1),
            ("The ‘lost opportunity’ to take the consulting job", 2),
            ([("Economic profits of flipping the house: ", {}),
              ("-50k (loss)", {'bold': True, 'color': RED})], 0, {}),
            ("= 700k − 600k − 150k", 1),
        ],
        size=24, sub_size=22, line_spacing_pts=12,
        notes=(
            "Adopted from Melanie Wasserman's Module 1 deck. Walk through "
            "the ledger: revenues of 700k, explicit costs of 600k (the "
            "purchase plus the remodel), and the implicit cost of 150k — "
            "the consulting salary you give up. The economic profit is "
            "negative 50k, so flipping the house is a loss once the "
            "opportunity cost is counted. The next slide adds one more "
            "possible opportunity cost."),
    )


def slide_44_another_opp(prs):
    def extras(slide):
        pic = _add_media_image(slide, "ic_s40_rId3.png",
                               left=Inches(2.95), top=Inches(5.15),
                               width=Inches(7.45), rounded=False,
                               shadow=True)
        pic.click_action.hyperlink.address = (
            "https://www.nytimes.com/2023/12/01/podcasts/the-daily/"
            "should-you-rent-or-buy-the")

    slide = content_slide(
        prs, 47, TAG_OPP, "Another Possible Opportunity Cost",
        [
            ("Additionally:", 0),
            ("What about the opportunity cost from the 500K of investment? "
             "That could also add to opportunity cost.", 1),
            ([("Podcast", {'underline': True}),
              (" on opportunity costs of buying a house (to be discussed "
               "during Thursday’s Coffee&Econ:", {})], 1, {}),
        ],
        size=26, sub_size=24, line_spacing_pts=14,
        bullets_top=Inches(1.7), bullets_height=Inches(3.1),
        extras=extras,
    )
    _link_runs(slide, {"Podcast": "https://www.nytimes.com/2023/12/01/"
                                  "podcasts/the-daily/should-you-rent-or-"
                                  "buy-the"})
    return slide


def slide_45_child_cost(prs):
    def extras(slide):
        _add_media_image(slide, "ic_s41_rId3.png",
                         left=Inches(2.65), top=Inches(2.55),
                         width=Inches(8.05), rounded=False, shadow=True)

    return content_slide(
        prs, 48, TAG_OPP, "The (Full) Cost of Raising a Child",
        [
            ("Data from the U.S. Department of Agriculture: Costs of "
             "raising a child born in the US in 2013, to age 18:", 0),
            ("", 0), ("", 0), ("", 0), ("", 0), ("", 0), ("", 0), ("", 0),
            ("Which costs are missing?", 0),
            ("  ", 1), ("  ", 1),
        ],
        size=22, sub_size=20, line_spacing_pts=10,
        bullets_top=Inches(1.5), bullets_height=Inches(5.5),
        extras=extras,
        notes=(
            "Source: http://graphics.wsj.com/childcost/ (using data from "
            "U.S. Department of Agriculture)\n"
            "Missing costs: “\"The real costs of raising a child for a "
            "moderate-income family\"—including forgone income, college "
            "for those who attend, and the so-called opportunity cost of "
            "not investing the money—\"would be closer to $900,000 to age "
            "22…\" says John Ward, an economist and the president of John "
            "Ward Economics, based in Prairie Village, Kan., which "
            "consults on legal disputes for plaintiffs and defendants.”"),
    )


def slide_46_child_penalty(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_OPP)
    _add_media_image(slide, "ic_s42_rId3.png",
                     left=Inches(2.65), top=Inches(0.60),
                     width=Inches(8.05), rounded=False, shadow=True)
    _draw_footer(slide, FOOTER_TEXT, 49)
    _set_notes(slide, (
        "Note: “zero trend” for men here simply means that wages don’t "
        "change after child birth RELATIVE to the trend for people w/o "
        "children. So the actual wage may well go up, even for women with "
        "children in the long run."))
    return slide


def slide_47_us_2022(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_OPP)
    _draw_action_title(slide, "Similar Figures for the US, Estimated in 2022")
    _add_media_image(slide, "ic_s43_rId2.jpeg",
                     left=Inches(1.35), top=Inches(1.85),
                     width=Inches(4.7), rounded=False, shadow=True)
    _add_media_image(slide, "ic_s43_rId3.jpeg",
                     left=Inches(6.45), top=Inches(1.85),
                     width=Inches(5.65), rounded=False, shadow=True)
    _draw_footer(slide, FOOTER_TEXT, 50)
    return slide


# --------------------------------------------------------------------------
# Slides 48–52 — sunk costs (was #44–48)
# --------------------------------------------------------------------------

def slide_48_outline_sunk(prs):
    return make_m1_outline(prs, 51, highlight_idx=4, tag=TAG_SUNK)


def slide_49_sunk_costs(prs):
    return content_slide(
        prs, 52, TAG_SUNK, "Sunk Costs",
        [
            ([("Sunk cost: ", {'color': RED}),
              ("A cost that has been paid and cannot be recovered",
               {})], 0, {}),
            ([("Ignore sunk costs: ", {'color': RED}),
              ("If the cost is sunk, it is not relevant for economic "
               "decisions", {})], 0, {}),
        ],
        size=28, sub_size=24, line_spacing_pts=24,
    )


def slide_50_sunk_examples(prs):
    return content_slide(
        prs, 53, TAG_SUNK, "Examples of Sunk Costs (you should ignore!)",
        [
            ("Individual consumers", 0),
            ("Dessert in fixed-price dinners", 1),
            ("Annual gym memberships (after having paid)", 1),
            ("Ski tickets (after buying, return impossible)", 1),
            ("Firms", 0),
            ("License fees (after paying)", 1),
            ("Advertising (after paying)", 1),
            ("To be continued during Module 3!", 1),
        ],
        size=26, sub_size=24, line_spacing_pts=12,
    )


def slide_51_concorde(prs):
    def extras(slide):
        _add_media_image(slide, "ic_s47_rId2.jpg",
                         left=Inches(7.55), top=Inches(2.15),
                         width=Inches(5.05), rounded=True)

    return content_slide(
        prs, 54, TAG_SUNK, "The Case of Concorde",
        [
            ("Very coveted and salient innovation:", 0),
            ("New York to London in under 3 hours.", 1),
            ("Pride of Air France and British Airways", 1),
            ("Massive financial losses", 0),
            ("Never recovered investments from innovation", 1),
            ("Did not even cover operating costs", 1),
            ("Economically rational decision would have been: shut down "
             "as soon as possible.", 0),
        ],
        size=24, sub_size=22, line_spacing_pts=12,
        bullets_width=Inches(6.9),
        extras=extras,
    )


def slide_52_sunk_takeaway(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_SUNK)
    _draw_action_title(slide, "Sunk Costs: Take-Away")
    left_items = [
        ([("Make an optimal decision ", {}),
          ("looking forward", {'underline': True}),
          (".", {})], 0, {}),
        ("As you cannot change the past, optimize the present/future.",
         0),
    ]
    right_items = [
        ("Often sunk costs come from unforeseen circumstances that change "
         "the optimal choice.", 0),
        ("Re-optimize without considering the past costs, if they cannot "
         "be recovered.", 0),
    ]
    b1 = _add_hierarchical_bullets(
        slide, left=MARGIN + Inches(0.3), top=Inches(2.0),
        width=Inches(5.9), height=Inches(4.4),
        items=left_items, size=26, line_spacing_pts=36)
    b1.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    b2 = _add_hierarchical_bullets(
        slide, left=Inches(7.0), top=Inches(2.0),
        width=Inches(5.9), height=Inches(4.4),
        items=right_items, size=26, line_spacing_pts=36)
    b2.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    _draw_footer(slide, FOOTER_TEXT, 55)
    return slide


# --------------------------------------------------------------------------
# Slides 53–56 — cost-benefit and marginal analysis (was #49–52)
# --------------------------------------------------------------------------

def slide_53_outline_cba(prs):
    return make_m1_outline(prs, 56, highlight_idx=5, tag=TAG_CBA)


def slide_54_cba(prs):
    def extras(slide):
        pic = _add_media_image(slide, "ic_s50_rId3.png",
                               left=Inches(9.55), top=Inches(0.60),
                               width=Inches(3.00), rounded=False,
                               shadow=False)
        pic.click_action.hyperlink.address = (
            "https://podcasts.apple.com/us/podcast/core-principle-1-the-"
            "cost-benefit-principl")

    slide = content_slide(
        prs, 57, TAG_CBA, "Cost-Benefit and Marginal Analysis",
        [
            ([("Objective: ", {}),
              ("Maximize ", {'color': RED}),
              ("net", {'bold': True, 'color': RED}),
              (" benefits ", {'color': RED}),
              ("(= total benefits – total costs)", {})], 0, {}),
            ([("To maximize net benefits from any activity, ", {}),
              ("marginal analysis", {'color': RED}),
              (" is required", {})], 0, {}),
            ([("Marginal benefit (MB): ", {'color': RED}),
              ("additional benefit due to an extra unit of the activity "
               "(or an extra unit of a good/service)", {})], 1, {}),
            ([("Marginal cost (MC): ", {'color': RED}),
              ("additional cost due to an extra unit", {})], 1, {}),
            ("Includes opportunity costs!", 2),
            ("Buy / go for the extra unit if MB > MC", 0),
            ("Important Rule:", 0,
             {'italic': True, 'underline': True, 'bullet_style': 'none'}),
            ([("Net benefits are maximized", {'color': RED}),
              (" where ", {}),
              ("MB = MC", {'bold': True})], 0, {}),
        ],
        size=24, sub_size=22, line_spacing_pts=11,
        extras=extras,
    )
    _highlight_texts(slide, ["Includes opportunity costs!"])
    return slide


def slide_55_exercise(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_CBA)
    _draw_action_title(slide,
                       "Example: How Many Hours a Week Should You Exercise?")

    base_y = 5.75                       # bar baseline (inches)
    scale = 1.0                          # inches per value unit
    mb = [3.0, 2.55, 2.1, 1.65, 1.2]
    mc = [1.0, 1.25, 1.5, 1.65, 2.0]
    col_w = 0.95
    gap = 0.42
    x0 = 6.05

    # Left summary: total net benefit bar + red double arrow
    _add_text(slide, Inches(0.7), Inches(1.95), Inches(2.6), Inches(0.95),
              "", size=20, color=NAVY)   # spacer for grouping stability
    tb = slide.shapes.add_textbox(int(Inches(0.55)), int(Inches(1.95)),
                                  int(Inches(2.9)), int(Inches(1.0)))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r1 = p.add_run(); r1.text = "Total "
    r2 = p.add_run(); r2.text = "Net"
    r3 = p.add_run(); r3.text = " Benefit"
    for r in (r1, r2, r3):
        r.font.name = "Calibri"; r.font.size = Pt(22)
        r.font.color.rgb = NB_BLUE
    r2.font.bold = True
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r4 = p2.add_run(); r4.text = "(in blue)"
    r4.font.name = "Calibri"; r4.font.size = Pt(18)
    r4.font.italic = True; r4.font.color.rgb = NB_BLUE
    total_net = sum(max(b - c, 0) for b, c in zip(mb, mc))   # 3.85
    bar_h = total_net * 0.72
    _add_rect(slide, Inches(1.35), Inches(base_y - bar_h),
              Inches(0.95), Inches(bar_h), NB_BLUE)
    arr = _add_arrow(slide, (Inches(2.62), Inches(base_y)),
                     (Inches(2.62), Inches(base_y - bar_h)),
                     color=RED, weight_pt=3.0, head=True)
    ln = arr.line._get_or_add_ln()
    hd = ET.SubElement(ln, qn('a:headEnd'))
    hd.set('type', 'triangle'); hd.set('w', 'med'); hd.set('h', 'med')

    # "Optimum where MB=MC" caption top-right
    cap = slide.shapes.add_textbox(int(Inches(7.1)), int(Inches(1.62)),
                                   int(Inches(5.3)), int(Inches(0.45)))
    ctf = cap.text_frame
    cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
    c1 = cp.add_run(); c1.text = "Optimum where "
    c1.font.name = "Calibri"; c1.font.size = Pt(22); c1.font.color.rgb = NAVY
    c2 = cp.add_run(); c2.text = "MB=MC"
    c2.font.name = "Calibri"; c2.font.size = Pt(22)
    c2.font.bold = True; c2.font.color.rgb = RED

    for k in range(5):
        x = x0 + k * (col_w + gap)
        mbv, mcv = mb[k] * scale * 0.95, mc[k] * scale * 0.95
        # hour label above the column
        _add_text(slide, Inches(x - 0.12), Inches(base_y - mb[0] - 0.42),
                  Inches(col_w + 0.25), Inches(0.3), "Hour %d" % (k + 1),
                  size=16, bold=True, color=NAVY, font="Calibri",
                  align=PP_ALIGN.CENTER)
        if mbv > mcv:
            # blue net-benefit block sits on top of the red MC block
            _add_rect(slide, Inches(x), Inches(base_y - mbv),
                      Inches(col_w), Inches(mbv - mcv), NB_BLUE)
            _add_rect(slide, Inches(x), Inches(base_y - mcv),
                      Inches(col_w), Inches(mcv), DARKRED)
        elif abs(mbv - mcv) < 1e-6:
            _add_rect(slide, Inches(x), Inches(base_y - mcv),
                      Inches(col_w), Inches(mcv), DARKRED)
        else:
            # MC exceeds MB: red up to MB, darker cap above for the excess
            _add_rect(slide, Inches(x), Inches(base_y - mbv),
                      Inches(col_w), Inches(mbv), DARKRED)
            _add_rect(slide, Inches(x), Inches(base_y - mcv),
                      Inches(col_w), Inches(mcv - mbv),
                      RGBColor(0x6E, 0x0A, 0x1A))
        # white MC label inside the red base
        _add_text(slide, Inches(x), Inches(base_y - 0.42),
                  Inches(col_w), Inches(0.3), "MC", size=14, bold=True,
                  color=WHITE, font="Calibri", align=PP_ALIGN.CENTER)
        # net-benefit callouts for the first two hours
        if k < 2:
            mid_y = base_y - mcv - (mbv - mcv) / 2
            _add_text(slide, Inches(x), Inches(mid_y - 0.30),
                      Inches(col_w), Inches(0.62),
                      "Net Benefit of Hour %d" % (k + 1), size=9,
                      bold=True, color=WHITE, font="Calibri",
                      align=PP_ALIGN.CENTER)
        # verdict row under the columns
        if k < 3:
            _add_arrow_shape(slide, Inches(x + 0.08), Inches(base_y + 0.14),
                             Inches(col_w - 0.16), Inches(0.34),
                             direction="right", fill=NB_BLUE)
        elif k == 3:
            _add_text(slide, Inches(x - 0.22), Inches(base_y + 0.14),
                      Inches(col_w + 0.44), Inches(0.3), "indifferent",
                      size=13, italic=True, color=NAVY, font="Calibri",
                      align=PP_ALIGN.CENTER)
        else:
            _add_arrow_shape(slide, Inches(x + 0.08), Inches(base_y + 0.14),
                             Inches(col_w - 0.16), Inches(0.34),
                             direction="left", fill=RED)
            _add_media_image(slide, "ic_s51_rId2.png",
                             left=Inches(x + 0.06), top=Inches(base_y + 0.58),
                             width=Inches(0.82), rounded=False, shadow=False)

    _draw_footer(slide, FOOTER_TEXT, 58)
    return slide


def slide_56_continuous(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_CBA)
    _add_text(slide, MARGIN, Inches(0.55), RULE_W, Inches(0.7),
              "Maximizing Net Benefit: The Standard Analysis that You’ll "
              "See throughout the Class",
              size=26, bold=True, color=NAVY, font="Calibri")
    _add_rect(slide, MARGIN, Inches(1.25), RULE_W, Inches(0.02), RULE)
    _add_rect(slide, MARGIN, Inches(1.235), GOLD_W, Inches(0.05), GOLD)

    _add_text(slide, MARGIN + Inches(0.3), Inches(1.42), Inches(8.0),
              Inches(0.42), "We’ll generally consider the continuous case:",
              size=22, color=NAVY, font="Calibri")

    fig = SimpleFig(4.1, 6.40, 5.4, 3.5, 10, 10)
    _fig_axes(slide, fig, y_title="$ per unit",
              x_title="")
    _add_text(slide, Inches(fig.l + fig.w - 0.35), Inches(fig.b + 0.10),
              Inches(2.6), Inches(0.62), "Quantity\n(of good or activity)",
              size=16, bold=True, italic=True, color=NAVY, font="Calibri")
    # MB: y = 9 - x ; MC: y = 1 + 0.75x ; intersection (4.57, 4.43)
    _fig_line(slide, fig, (0.5, 8.5), (7.5, 1.5), color=GREEN_DK,
              weight_pt=2.75)
    _fig_line(slide, fig, (0.5, 1.375), (8.5, 7.375), color=RED,
              weight_pt=2.75)
    _fig_guide(slide, fig, (4.5714, 4.4286), to_y=False, color=NAVY,
               dash='sysDot')
    _add_text(slide, fig.x(4.5714) - Inches(0.85), Inches(fig.b + 0.08),
              Inches(1.9), Inches(0.32), "Q* (optimum)", size=18,
              bold=True, italic=True, color=NAVY, font="Calibri",
              align=PP_ALIGN.CENTER)
    _add_text(slide, fig.x(7.6), fig.y(1.7), Inches(2.6), Inches(0.35),
              "MB (Marginal Benefit)", size=18, bold=True, italic=True,
              color=GREEN_DK, font="Calibri")
    _add_text(slide, fig.x(6.9), fig.y(7.7), Inches(2.6), Inches(0.35),
              "MC (Marginal Cost)", size=18, bold=True, italic=True,
              color=RED, font="Calibri")
    _draw_footer(slide, FOOTER_TEXT, 59)
    return slide


# --------------------------------------------------------------------------
# Slides 57–58 — summary + next steps (was #53 + MW #71)
# --------------------------------------------------------------------------

def slide_57_summary(prs):
    return content_slide(
        prs, 60, TAG_WRAP, "Module 1 Summary",
        [
            ("Markets, supply, and demand", 0),
            ("Market equilibrium: intersection of supply and demand", 0),
            ("Key economic principles for decision-making", 0),
            ("Economic costs include opportunity costs", 1),
            ([("Economists count ", {}),
              ("implicit costs", {'italic': True, 'underline': True}),
              (" as well (in contrast to accountants).", {})], 2, {}),
            ("Ignore sunk costs", 1),
            ("Use cost-benefit and marginal analysis: MB=MC", 1),
        ],
        size=26, sub_size=24, line_spacing_pts=14,
    )


def slide_58_next_steps(prs):
    return content_slide(
        prs, 61, TAG_WRAP, "Next Steps",
        [
            ("Problem set 1 posted", 0),
            ("Due [DATE]", 1),
            ("TA will send submission instructions via email/BruinLearn", 1),
            ("Optional practice exercises in the textbook and Achieve", 0),
            ("Module 2: Demand Analysis", 0),
            ("Textbook reading", 1),
            ("Read news article I will send via email/BruinLearn", 1),
            ("Take survey I will send via email/BruinLearn", 1),
        ],
        size=26, sub_size=24, line_spacing_pts=14,
        notes=(
            "Adopted from Melanie Wasserman's Module 1 deck as a closing "
            "slide: problem set logistics plus a preview of Module 2. The "
            "due date is a placeholder; the last two sub-bullets are her "
            "pre-class routine — keep or cut as fits Nico's flow."),
    )



# ==========================================================================
#  BATCH D — Videos 1–4 (slides 59–84) + build orchestration
# ==========================================================================

def make_video_title(prs, topic, video_no):
    """Video title slide in the deck-title format (self-contained for the
    later split into separate video decks)."""
    slide = _blank_slide(prs)
    _add_text(slide, Inches(0.9), Inches(2.05), SLIDE_W - Inches(1.8),
              Inches(1.1), topic,
              size=54, bold=True, color=NAVY, font="Calibri",
              align=PP_ALIGN.CENTER)
    _add_text(slide, 0, Inches(3.45), SLIDE_W, Inches(0.75),
              "Module 1 – Video %d" % video_no,
              size=40, bold=True, color=GOLD, font="Calibri",
              align=PP_ALIGN.CENTER)
    _add_rect(slide, int((SLIDE_W - Inches(4.0)) / 2), Inches(4.48),
              Inches(4.0), 54864, GOLD)
    _add_text(slide, 0, Inches(4.85), SLIDE_W, Inches(0.55),
              "Management 405", size=26, bold=True, color=GRAY,
              font="Calibri", align=PP_ALIGN.CENTER)
    _add_text(slide, 0, Inches(5.50), SLIDE_W, Inches(0.5),
              "Prof. Nico Voigtländer  ·  UCLA Anderson",
              size=22, color=GRAY, font="Calibri", align=PP_ALIGN.CENTER)
    _add_rect(slide, 0, Inches(7.15), SLIDE_W, Inches(0.02), RULE)
    _add_rect(slide, MARGIN, Inches(7.135), GOLD_W, Inches(0.05), GOLD)
    return slide


# ---- Video 1 (59–61) ------------------------------------------------------

def slide_59_v1_title(prs):
    return make_video_title(prs, "Introduction", 1)


def slide_60_v1_roadmap(prs):
    return make_roadmap(prs, 63, tag=TAG_V1)


def slide_61_v1_outline(prs):
    return make_m1_outline(prs, 64, tag=TAG_V1, descriptions=True)


# ---- Video 2 (62–66) ------------------------------------------------------

def slide_62_v2_title(prs):
    return make_video_title(prs, "Markets", 2)


def slide_63_v2_outline(prs):
    return make_m1_outline(prs, 66, tag=TAG_V2, highlight_idx=0)


def slide_64_v2_market_def(prs):
    return content_slide(
        prs, 67, TAG_V2, "Market Definition",
        [
            ("A company must understand its market", 0,
             {'bold': True, 'bullet_style': 'none'}),
            ("Customers", 0),
            ("Competitors (actual and potential)", 0),
            ("Extent of market", 0,
             {'bold': True, 'bullet_style': 'none'}),
            ("Which products belong to a market?", 0),
            ("Simple test to identify the range of products in your "
             "market: If the price of another product changes, will demand "
             "for your product change?", 1),
            ("Relevant for antitrust litigation (in mergers & "
             "acquisitions)", 1),
            ("Geography boundaries", 0),
            ("Coffee shop in Venice (CA) v. gasoline retail v. gold", 1),
        ],
        size=24, sub_size=22, line_spacing_pts=10,
    )


# --------------------------------------------------------------------------
# Displays 73-74 — Tapestry-Capri market-definition mini-case
# --------------------------------------------------------------------------
# Nico copied these in from "Module 1 - Example Candidates.pptx" (slides 2
# and 3) on 2026-08-23 and animated them. Ported here so build.py stays the
# source of truth; his choreography lives in _animate.py PLANS 73/74.

def _photo_caption(slide, left, top, width,
                   text="Photos: Wikimedia Commons"):
    return _add_text(slide, left, top, width, Inches(0.28), text,
                     size=11, italic=True, color=GRAY, font="Calibri",
                     align=PP_ALIGN.CENTER)


def _quote_box(slide, left, top, width, height, quote, attribution, *,
               size=18):
    """Verbatim quote in the cream convention box (Teaching CLAUDE.md)."""
    return _add_convention_box(
        slide, left, top, width, height,
        runs=[(quote, {'italic': True, 'size': size}),
              ("   \u2014 " + attribution,
               {'bold': True, 'size': max(size - 2, 16),
                'newline': True})],
        align=PP_ALIGN.LEFT)


def slide_tapestry_case(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_V2)
    _draw_action_title(slide, "Market Definition Mini-Case: Tapestry\u2013Capri")
    box = _add_hierarchical_bullets(
        slide, left=MARGIN, top=Inches(1.50), width=Inches(8.0),
        height=Inches(5.3),
        items=[
            ("Aug 2023: Tapestry (Coach, Kate Spade) agrees to buy Capri "
             "(Michael Kors, Versace) for $8.5B", 0),
            ([("The definition of the market", {'bold': True}),
              (" would turn out to be crucial for the case:", {})], 0, {}),
            ([("FTC: the market is ", {}),
              ("\u201caccessible luxury\u201d handbags", {'color': RED}),
              (" \u2014 roughly $100 to under $1,000", {})], 1, {}),
            ([("The firms: the market is ", {}),
              ("all handbags", {'color': RED}),
              (" \u2014 from fast fashion to Herm\u00e8s \u2014 and entry "
               "is easy", {})], 1, {}),
        ],
        size=24, sub_size=22, line_spacing_pts=16,
        sub_line_spacing_pts=8)
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    _add_media_image(slide, "web_coach.jpg",
                     left=Inches(8.55), top=Inches(1.70),
                     width=Inches(3.67), rounded=True)
    _add_media_image(slide, "web_michaelkors.jpg",
                     left=Inches(8.75), top=Inches(4.30),
                     width=Inches(3.27), rounded=True)
    _photo_caption(slide, Inches(8.55), Inches(6.80), Inches(3.67))
    _draw_footer(slide, FOOTER_TEXT, 73)
    _set_notes(slide, (
        "In August 2023 Tapestry, which owns Coach and Kate Spade, agreed "
        "to buy Capri, which owns Michael Kors and Versace, for 8.5 "
        "billion dollars. The whole case then turned on one question: what "
        "is the market? The FTC said it was \u201caccessible luxury\u201d "
        "handbags, roughly 100 dollars to just under a thousand, where "
        "these brands sit almost on top of each other. The firms said the "
        "market was all handbags, from fast fashion to Herm\u00e8s, where "
        "their combined share looks small and entry is easy. Hold that "
        "disagreement in mind, because the next slide shows which side the "
        "companies' own documents supported.\n"
        "Case record: FTC administrative suit filed 22 April 2024; "
        "preliminary injunction granted by Judge Jennifer Rochon, "
        "S.D.N.Y., 24 October 2024; merger agreement terminated 13 "
        "November 2024. The defense called the FTC's market "
        "\u201cgerrymandered\u201d; the court found the accessible-luxury "
        "segment real, pointing to distinct prices, customers, "
        "discounting and outlet distribution. Photos: Wikimedia Commons "
        "(Coach store, Tenmaya Fukuyama; Michael Kors store, Rehoboth "
        "Beach DE)."))
    return slide


def slide_tapestry_evidence(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_V2)
    _draw_action_title(slide,
                       "The Firms' Own Documents Drew the Market Boundary")
    # price-tier ladder
    y0, h0 = Inches(1.55), Inches(1.30)
    _add_rounded_filled_box(slide, Inches(0.75), y0, Inches(3.55), h0,
                            "Mass market\nunder $100",
                            fill=FADED, text_color=WHITE, size=19)
    _add_rounded_filled_box(
        slide, Inches(4.55), y0, Inches(4.25), h0,
        "\u201cAccessible luxury\u201d\n$100 \u2013 under $1,000\n"
        "Coach \u00b7 Kate Spade \u00b7 Michael Kors",
        fill=GOLD, text_color=NAVY, size=18)
    _add_rounded_filled_box(slide, Inches(9.05), y0, Inches(3.55), h0,
                            "True luxury\n$1,000+ \u00b7 LV \u00b7 "
                            "Herm\u00e8s \u00b7 Chanel",
                            fill=FADED, text_color=WHITE, size=19)
    # (the gold "the FTC's market ..." line that sat here was deleted by
    #  Nico on 2026-08-23; the point is in the speaker notes)
    _add_text(slide, MARGIN, Inches(3.32), RULE_W, Inches(0.35),
              "Combined Tapestry + Capri share of \u201caccessible "
              "luxury\u201d handbags:", size=18, bold=True, color=NAVY,
              font="Calibri", align=PP_ALIGN.CENTER)
    cards = [("59%", "FTC's expert (third-party data)"),
             ("77%", "Capri's internal documents"),
             ("83%", "Tapestry's internal data")]
    for i, (num, lab) in enumerate(cards):
        x = Inches(1.35 + i * 3.65)
        shp = _add_outlined_box(slide, x, Inches(3.72), Inches(3.35),
                                Inches(1.10), "", rounded=True,
                                shadow=True, line=GOLD, line_w=1.5)
        tf = shp.text_frame
        para = tf.paragraphs[0]
        # _add_outlined_box seeds an empty run; drop it so the card has
        # exactly the two runs Nico's version has
        for stale in list(para.runs):
            stale._r.getparent().remove(stale._r)
        r1 = para.add_run(); r1.text = num + "   "
        r1.font.name = "Calibri"; r1.font.size = Pt(30)
        r1.font.bold = True; r1.font.color.rgb = NAVY
        r2 = para.add_run(); r2.text = lab
        r2.font.name = "Calibri"; r2.font.size = Pt(18)
        r2.font.color.rgb = GRAY
    _add_text(slide, MARGIN, Inches(4.92), RULE_W, Inches(0.30),
              "(figures from documents the companies had to hand over in "
              "the merger review \u2014 not leaked)", size=18, italic=True,
              color=GRAY, font="Calibri", align=PP_ALIGN.CENTER)
    _quote_box(slide, MARGIN + Inches(0.35), Inches(5.34),
               RULE_W - Inches(0.7), Inches(1.00),
               "\u201cBottom line, saying we're in the same market with "
               "true luxury is a joke. \u2026 Nobody says \u2018should I "
               "buy a LV bag or a Coach bag?\u2019\u201d",
               "internal Tapestry message cited by the court")
    _add_rounded_filled_box(
        slide, Inches(0.90), Inches(6.44), Inches(11.53), Inches(0.60),
        "Oct 2024: the court sides with the FTC and blocks the deal "
        "\u2014 merger abandoned Nov 2024",
        fill=GOLD, text_color=NAVY, size=19, bold=True, corner_pct=0.18)
    _draw_footer(slide, FOOTER_TEXT, 74)
    _set_notes(slide, (
        "This is where the case was decided. The three cards are the same "
        "quantity \u2014 the combined share of accessible-luxury handbags "
        "\u2014 computed three ways: about 59 percent by the FTC's "
        "economic expert using largely third-party data, 77 percent in "
        "Capri's own internal documents, and 83 percent from Tapestry's "
        "own internal data. Note where those last two came from: "
        "ordinary-course documents the companies had to produce in the "
        "merger review, not leaks. And the internal message at the bottom "
        "is the firms' own people rejecting the very market definition "
        "their lawyers were arguing for. In October 2024 the court sided "
        "with the FTC and blocked the deal, and the merger was abandoned "
        "the following month.\n"
        "One more detail worth knowing: \u201caccessible luxury\u201d was "
        "the firms' own term, used extensively in SEC filings and investor "
        "presentations, and it disappeared from their vocabulary once the "
        "FTC sued. Sources: Judge Rochon's opinion, S.D.N.Y., 24 October "
        "2024 (share figures at 97); Clifford Chance and MoFo case "
        "notes."))
    return slide


def slide_65_v2_netflix(prs):
    def extras(slide):
        _add_media_image(slide, "v2_s04_rId4.png",
                         left=Inches(7.85), top=Inches(2.45),
                         width=Inches(4.85), rounded=True)
        _add_discussion_break(slide, text="Prepare for Class Discussion",
                              width=Inches(5.9))

    return content_slide(
        prs, 68, TAG_V2, "Market Definition: the Case of Netflix",
        [
            ("Define Netflix’s market", 0,
             {'bold': True, 'bullet_style': 'none'}),
            ("Online streaming?", 0),
            ("All films?", 0),
            ("All entertainment?", 0),
            ("Did the market change with Covid-19?", 0),
        ],
        size=28, sub_size=24, line_spacing_pts=16,
        bullets_width=Inches(7.0), bullets_height=Inches(4.2),
        extras=extras,
    )


def slide_66_v2_actors(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_V2)
    _draw_action_title(slide, "Actors in the Market")

    cards = [
        ("CONSUMERS", [
            ("Objective: Maximize “utility” from goods and services, "
             "given their limited income", 0),
        ]),
        ("WORKERS", [
            ("Choose a job to earn income", 0),
            ("Face trade-offs in their choice of employment", 0),
            ("E.g. labor income vs. leisure time / flexibility", 1),
        ]),
        ("FIRMS", [
            ("Employ workers to produce goods and services for "
             "consumers", 0),
            ("Objective: Maximize profits (value of the firm)", 0),
        ]),
    ]
    card_w = Inches(4.0)
    gap = Inches(0.25)
    x0 = (SLIDE_W - 3 * card_w - 2 * gap) // 2
    for i, (header, items) in enumerate(cards):
        x = x0 + i * (card_w + gap)
        _add_rounded_filled_box(slide, x, Inches(1.95), card_w,
                                Inches(0.62), header,
                                fill=NAVY, text_color=WHITE, size=22,
                                bold=True)
        box = _add_hierarchical_bullets(
            slide, left=x + Inches(0.12), top=Inches(2.80),
            width=card_w - Inches(0.24), height=Inches(3.6),
            items=items, size=20, sub_size=18, line_spacing_pts=10)
        box.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    _draw_footer(slide, FOOTER_TEXT, 69)
    return slide


# ---- Video 3 (67–76) ------------------------------------------------------

def slide_67_v3_title(prs):
    return make_video_title(prs, "Supply and Demand", 3)


def slide_68_v3_outline(prs):
    return make_m1_outline(prs, 71, tag=TAG_V3, highlight_idx=1)


def slide_69_v3_ds_analysis(prs):
    return content_slide(
        prs, 72, TAG_V3, "Demand-Supply Analysis",
        [
            ("How do supply and demand determine market prices and the "
             "quantity produced/consumed?", 0),
            ("Wide variety of applications", 0),
            ("Example:", 0),
            ("How does growth in China affect the price of oil? The price "
             "of textiles? The price of (electric) cars?", 1),
        ],
        size=28, sub_size=24, line_spacing_pts=18,
    )


DEMAND_DEF_NOTE = (
    "Note: aside from price, we are holding everything else constant, "
    "including consumers’ income, tastes, price of other goods.\n"
    "Thought experiment: suppose you show up at Pike Place Market on a "
    "Saturday and observe the price of tomatoes and total quantity that "
    "consumers buy. Suppose you can go back in time to that same Saturday "
    "and replace all the price tags using an alternative price and see the "
    "quantity that consumers buy. Do this over and over until you trace "
    "out the demand curve. If tomatoes are $2/pound in one iteration, some "
    "people will buy. If they are $1/pound, additional people will buy. As "
    "price declines, more consumers want to purchase it. This will make "
    "the demand curve downward sloping.\n"
    "We will get into this in much more detail next week.")


def _add_definition_box(slide, body, *, top=Inches(1.70),
                        left=Inches(1.35), width=Inches(10.65),
                        height=Inches(1.15), size=20):
    """Cream definition callout: bold 'Definition:' prefix + body with an
    italic-underlined 'holding everything else constant' tail."""
    runs = [("Definition", {'bold': True, 'size': size}),
            (": " + body, {'size': size}),
            ("holding everything else constant",
             {'italic': True, 'underline': True, 'size': size})]
    return _add_convention_box(slide, left, top, width, height, runs=runs)


def slide_70_v3_demand_def(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_V3)
    _draw_action_title(slide, "The Market Demand Curve")
    _add_definition_box(
        slide,
        "the relationship between the quantity of a good that consumers "
        "are willing to buy and the price of the good, ",
        # hand-tweaked 2026-08-23 from left 1.35" / top 3.35"
        left=1417320, top=1422849, height=Inches(1.5), size=24)
    _draw_footer(slide, FOOTER_TEXT, 73)
    _set_notes(slide, DEMAND_DEF_NOTE)
    return slide


def slide_71_v3_ceteris_paribus(prs):
    def extras(slide):
        # named so _group_pass.py can merge the header with the picture
        # (Nico grouped them by hand on 2026-08-23)
        hdr = _add_text(slide, Inches(8.05), Inches(1.95), Inches(4.6),
                        Inches(0.62),
                        "Things are not held constant between these two "
                        "cones",
                        size=16, bold=True, color=NAVY, font="Calibri",
                        align=PP_ALIGN.CENTER)
        hdr.name = "sdlabel:cones"
        pic = _add_media_image(slide, "v3_s05_rId3.png",
                               left=Inches(8.25), top=Inches(2.60),
                               width=Inches(4.25), rounded=True)
        pic.name = "sdpic:cones"

    slide = content_slide(
        prs, 74, TAG_V3,
        "“Holding Everything Else Constant” – “Ceteris Paribus” (c.p.)",
        [
            ([("Want to know:", {'color': RED}),
              (" how the quantity demanded of an ice cream cone changes "
               "with a price increase", {})], 0, {}),
            ([("Assume: ", {'color': RED}),
              ("Everything else that affects the quantity demanded "
               "remains constant", {})], 0, {}),
            ("What is held constant?", 0, {'color': RED}),
            ("Consumer income", 1),
            ("Outside temperature", 1),
            ("Prices of related goods", 1),
            ("Size, shape, quality of the ice cream cone", 1),
            ([("…everything", {'italic': True}),
              (" other than price", {})], 1, {}),
        ],
        size=22, sub_size=20, line_spacing_pts=10,
        bullets_width=Inches(7.3), title_size=26,
        extras=extras,
        notes=DEMAND_DEF_NOTE,
    )
    return slide


def slide_72_v3_demand_curve(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_V3)
    _draw_action_title(slide, "The Market Demand Curve")
    _add_definition_box(
        slide,
        "the relationship between the quantity of a good that consumers "
        "are willing to buy and the price of the good, ",
        size=18, height=Inches(0.95))
    box = _add_hierarchical_bullets(
        slide, left=MARGIN + Inches(0.2), top=Inches(3.15),
        width=Inches(6.6), height=Inches(3.2),
        items=[
            ([("The demand curve is ", {}),
              ("downward sloping:", {'italic': True, 'color': RED}),
              (" consumers want to purchase more of a good as its price "
               "goes down, ", {}),
              ("all else constant", {'italic': True})], 0,
             {'bullet_style': 'none'}),
        ],
        size=24, line_spacing_pts=0)
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    fig = SimpleFig(8.6, 6.55, 3.9, 3.0, 10, 10)
    _fig_axes(slide, fig, label_size=16)
    _fig_line(slide, fig, (0.8, 8.6), (7.6, 1.4), color=GOLD,
              weight_pt=2.75)
    _fig_curve_label(slide, fig, 7.8, 1.9, "D", color=NAVY)
    _draw_footer(slide, FOOTER_TEXT, 75)
    _set_notes(slide, DEMAND_DEF_NOTE)
    return slide


def slide_73_v3_move_vs_shift_d(prs):
    def extras(slide):
        fig = SimpleFig(8.1, 6.40, 4.3, 4.05, 10, 10)
        _fig_axes(slide, fig, label_size=16)
        # D: y = 9 - x ; movement (4.5,4.5) -> (2.5,6.5); D': y = 12 - x
        _fig_guide(slide, fig, (4.5, 4.5), color=GRAY)
        _fig_guide(slide, fig, (2.5, 6.5), color=GRAY)
        _fig_line(slide, fig, (1, 8), (8, 1), color=GOLD, weight_pt=2.75)
        _fig_curve_label(slide, fig, 8.05, 1.35, "D", color=GOLD)
        _fig_line(slide, fig, (2.5, 6.5), (4.5, 4.5), color=RED,
                  weight_pt=4.5)
        _add_text(slide, fig.x(2.55), fig.y(7.6), Inches(0.6), Inches(0.3),
                  "i)", size=18, bold=True, italic=True, color=RED,
                  font="Calibri")
        # 2026-08-23 (Nico): darker green throughout, D’ label / ii)
        # arrow / ii) label repositioned by hand, and a new horizontal
        # dashed segment carrying P1 out to Q3. The D’ set and the Q3 set
        # are each grouped (names below drive _group_pass.py rule 5).
        dprime = _fig_line(slide, fig, (3.6, 8.4), (9.7, 2.3),
                           color=GREEN_DK, weight_pt=2.75, dash='dash')
        dprime.name = "sdcurve:Dp"
        dplab = _fig_curve_label(slide, fig, 9.6070, 2.9284, "D’",
                                 color=GREEN_DK)
        dplab.name = "sdlabel:Dp"
        _, q3v = _fig_guide(slide, fig, (7.5, 4.5), to_y=False, color=GRAY)
        if q3v is not None:
            q3v.name = "sdguide:v:Q3"
        # the new segment: P1 level, from Q1 across to Q3
        q3h = _add_arrow(slide, (fig.x(4.5), fig.y(4.5)),
                         (fig.x(7.5), fig.y(4.5)),
                         color=GRAY, weight_pt=1.25, head=False,
                         dash='dash')
        q3h.name = "sdguide:h:Q3"
        arr = _add_arrow(slide, (fig.x(4.7488), fig.y(4.6852)),
                         (fig.x(5.8488), fig.y(5.3852)),
                         color=GREEN_DK, weight_pt=2.0, head=True)
        arr.name = "sdarrow:ii"
        iilab = _add_text(slide, fig.x(4.8488), fig.y(5.8716), Inches(0.7),
                          Inches(0.3),
                          "ii)", size=18, bold=True, italic=True,
                          color=GREEN_DK, font="Calibri")
        iilab.name = "sdlabel:ii"
        _fig_ylab(slide, fig, 4.5, "P1", size=16)
        _fig_ylab(slide, fig, 6.5, "P2", size=16)
        _fig_xlab(slide, fig, 4.5, "Q1", size=16)
        _fig_xlab(slide, fig, 2.5, "Q2", size=16)
        _fig_xlab(slide, fig, 7.5, "Q3", size=16).name = "sdxlab:Q3"

    return content_slide(
        prs, 76, TAG_V3, "The Market Demand Curve",
        [
            ("Distinguish between:", 0,
             {'bold': True, 'bullet_style': 'none'}),
            ([("i) Movement along D:", {'italic': True, 'color': RED}),
              (" changes in quantity demanded when price changes, ", {}),
              ("all else constant", {'italic': True})], 0, {}),
            ([("ii) Shift in Demand:",
               {'italic': True, 'color': GREEN_DK}),
              (" changes in demand due to non-price factors", {})], 0, {}),
            ([("e.g. higher income shifts the demand curve to the right "
               "(from ", {}),
              ("D", {'italic': True}), (" to ", {}),
              ("D′", {'italic': True}), (")", {})], 1, {}),
        ],
        size=24, sub_size=22, line_spacing_pts=16,
        bullets_width=Inches(6.6),
        extras=extras,
    )


def slide_74_v3_ai_chips(prs):
    """Replaces the COVID/flour example (2026-08-23, Nico): the AI build-out
    as a non-price factor shifting the demand for chips outward. Deliberately
    no supply curve — Video 3 has not reached market equilibrium yet."""
    def extras(slide):
        # Nico's own image, dropped into _source_images on 2026-08-23
        # (converted from .webp — python-pptx cannot embed webp). No
        # source credit: it is not a Wikimedia Commons photo.
        pic = _add_media_image(slide, "AI_Accelerator_Chips.png",
                               left=Inches(0.95), top=Inches(3.70),
                               width=Inches(5.30), rounded=True)
        pic.name = "sdpic:chips"

        fig = SimpleFig(8.1, 6.40, 4.3, 4.05, 10, 10)
        _fig_axes(slide, fig, label_size=16)
        # D: y = 9 - x ; D': y = 12 - x (a parallel outward shift), so at
        # the price P1 = 4.5 the quantity demanded goes from 4.5 to 7.5
        gh, gv = _fig_guide(slide, fig, (4.5, 4.5), color=GRAY)
        if gh is not None:
            gh.name = "sdguide:h:P1"
        if gv is not None:
            gv.name = "sdguide:v:Q1"
        d = _fig_line(slide, fig, (1, 8), (8, 1), color=GOLD,
                      weight_pt=2.75)
        d.name = "sdcurve:D"
        _fig_curve_label(slide, fig, 8.05, 1.35, "D",
                         color=GOLD).name = "sdlabel:D"
        dp = _fig_line(slide, fig, (3.6, 8.4), (9.7, 2.3), color=GREEN_DK,
                       weight_pt=2.75, dash='dash')
        dp.name = "sdcurve:Dp"
        _fig_curve_label(slide, fig, 9.6070, 2.9284, "D’",
                         color=GREEN_DK).name = "sdlabel:Dp"
        # 2026-08-23 (Nico): no Q2 guides here. Holding the price fixed
        # while the quantity moves out would suggest the price does not
        # change, and this slide has no supply curve to settle that.
        arr = _add_arrow(slide, (fig.x(4.7488), fig.y(4.6852)),
                         (fig.x(5.8488), fig.y(5.3852)),
                         color=GREEN_DK, weight_pt=2.0, head=True)
        arr.name = "sdarrow:shift"
        _fig_ylab(slide, fig, 4.5, "P1", size=16).name = "sdylab:P1"
        _fig_xlab(slide, fig, 4.5, "Q1", size=16).name = "sdxlab:Q1"

    return content_slide(
        prs, 77, TAG_V3, "AI and the Demand for Computer Chips",
        [
            ("Firms are racing to build the capacity to train and run AI "
             "models", 0),
            ("Demand for AI accelerator chips is rising dramatically",
             0),
        ],
        size=24, sub_size=22, line_spacing_pts=16,
        bullets_width=Inches(6.6), bullets_top=Inches(1.60),
        bullets_height=Inches(1.85),
        extras=extras,
        notes=(
            "The AI build-out is a textbook non-price factor. Nothing about "
            "the price of a chip changed; what changed is how badly firms "
            "want them, because training and running AI models takes large "
            "numbers of specialised processors. So the whole demand curve "
            "moves out, from D to D prime. Read it at a fixed price: at P1 "
            "the quantity demanded goes from Q1 to Q2. Note that there is "
            "no supply curve here yet — we are still looking at demand on "
            "its own, and we will put the two sides together in Video 4. "
            "That is also why the chart does not mark a new quantity: "
            "without supply we cannot say what happens to the price."),
    )


SUPPLY_DEF_NOTE = (
    "As price rises, quantity supplied will increase. Intuition for upward "
    "sloping: many firms experience increasing costs of production as "
    "their output rises; need to earn higher price in order to be willing "
    "to sell more output. We will get into this in much more detail in "
    "Module 4.")


def slide_75_v3_supply_curve(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_V3)
    _draw_action_title(slide, "The Market Supply Curve")
    _add_convention_box(
        slide, Inches(1.35), Inches(1.70), Inches(10.65), Inches(0.95),
        runs=[("Definition", {'bold': True, 'size': 18}),
              (": Describes the relationship between the quantity of a "
               "good that producers are willing to sell and the price of "
               "the good, ", {'size': 18}),
              ("holding everything else constant",
               {'italic': True, 'underline': True, 'size': 18})])
    box = _add_hierarchical_bullets(
        slide, left=MARGIN + Inches(0.2), top=Inches(3.15),
        width=Inches(6.6), height=Inches(3.2),
        items=[
            ([("Upward sloping:", {'italic': True, 'color': RED}),
              (" the higher the price, the more firms are able and "
               "willing to produce and sell, ", {}),
              ("all else constant", {'italic': True})], 0,
             {'bullet_style': 'none'}),
        ],
        size=24, line_spacing_pts=0)
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    fig = SimpleFig(8.6, 6.55, 3.9, 3.0, 10, 10)
    _fig_axes(slide, fig, label_size=16)
    _fig_line(slide, fig, (0.8, 1.6), (7.6, 8.6), color=STEEL,
              weight_pt=2.75)
    _fig_curve_label(slide, fig, 7.8, 9.0, "S", color=NAVY)
    _draw_footer(slide, FOOTER_TEXT, 78)
    _set_notes(slide, SUPPLY_DEF_NOTE)
    return slide


def slide_76_v3_move_vs_shift_s(prs):
    def extras(slide):
        fig = SimpleFig(8.1, 6.40, 4.3, 4.05, 10, 10)
        _fig_axes(slide, fig, label_size=16)
        # S: y = x + 1 ; movement (4.5,5.5) -> (3,4); S': y = x - 1.5
        h2, v2 = _fig_guide(slide, fig, (4.5, 5.5), color=GRAY)
        if h2 is not None:
            h2.name = "sdguide:h:P2"
        if v2 is not None:
            v2.name = "sdguide:v:Q2"
        h1, v1 = _fig_guide(slide, fig, (3, 4), color=GRAY)
        if h1 is not None:
            h1.name = "sdguide:h:P1"
        if v1 is not None:
            v1.name = "sdguide:v:Q1"
        _fig_line(slide, fig, (1, 2), (7.5, 8.5), color=STEEL,
                  weight_pt=2.75)
        _fig_curve_label(slide, fig, 7.7, 8.9, "S", color=NAVY)
        _fig_line(slide, fig, (3, 4), (4.5, 5.5), color=RED,
                  weight_pt=4.5).name = "sdarrow:i"
        # hand-tweaked 2026-08-23 (from 2.9, 5.6)
        i_lab = _add_text(slide, fig.x(3.0512), fig.y(5.1481), Inches(0.6),
                          Inches(0.3),
                          "i)", size=18, bold=True, italic=True, color=RED,
                          font="Calibri")
        i_lab.name = "sdlabel:i"
        _fig_line(slide, fig, (2.5, 1), (9.6, 8.1), color=GREEN_DK,
                  weight_pt=2.75, dash='dash').name = "sdcurve:Sp"
        _fig_curve_label(slide, fig, 9.4, 8.6, "S’",
                         color=GREEN_DK).name = "sdlabel:Sp"
        _, q3v = _fig_guide(slide, fig, (7, 5.5), to_y=False, color=GRAY)
        if q3v is not None:
            q3v.name = "sdguide:v:Q3"
        # new 2026-08-23: carries P2 across from Q2 to Q3
        q3h = _add_arrow(slide, (fig.x(4.5), fig.y(5.5)),
                         (fig.x(7.0), fig.y(5.5)),
                         color=GRAY, weight_pt=1.25, head=False,
                         dash='dash')
        q3h.name = "sdguide:h:Q3"
        # arrow + label hand-tweaked 2026-08-23 (from 5.6->6.8 at 6.4,
        # and label at 6.0, 7.2)
        _add_arrow(slide, (fig.x(5.7860), fig.y(6.0815)),
                   (fig.x(6.9860), fig.y(6.0815)),
                   color=GREEN_DK, weight_pt=2.0,
                   head=True).name = "sdarrow:ii"
        _add_text(slide, fig.x(6.1860), fig.y(6.8815), Inches(0.7),
                  Inches(0.3),
                  "ii)", size=18, bold=True, italic=True, color=GREEN_DK,
                  font="Calibri").name = "sdlabel:ii"
        # 2026-08-23 (Nico): the movement along S runs from the LOWER
        # price to the higher one, so (4.0, 3.0) is point 1 and
        # (5.5, 4.5) is point 2 — the labels used to be the other way up.
        _fig_ylab(slide, fig, 5.5, "P2", size=16).name = "sdylab:P2"
        _fig_ylab(slide, fig, 4.0, "P1", size=16).name = "sdylab:P1"
        _fig_xlab(slide, fig, 4.5, "Q2", size=16).name = "sdxlab:Q2"
        _fig_xlab(slide, fig, 3.0, "Q1", size=16).name = "sdxlab:Q1"
        _fig_xlab(slide, fig, 7.0, "Q3", size=16).name = "sdxlab:Q3"

    return content_slide(
        prs, 79, TAG_V3, "The Market Supply Curve",
        [
            ("Distinguish between:", 0,
             {'bold': True, 'bullet_style': 'none'}),
            ([("i) Movement along S:", {'italic': True, 'color': RED}),
              (" changes in quantity supplied when price changes, ", {}),
              ("all else constant", {'italic': True})], 0, {}),
            # 2026-08-23 (Nico): sub-bullet under i), same format as the
            # sub-bullet under ii)
            ("Existing firms supply more when prices rise", 1),
            ([("ii) Shifts in S:", {'italic': True, 'color': GREEN_DK}),
              (" changes in supply; e.g. a shift to the right (from ", {}),
              ("S", {'italic': True}), (" to ", {}),
              ("S′", {'italic': True}), (")", {})], 0, {}),
            ("Example: better technology, lower input prices, new firms "
             "enter the market", 1),
        ],
        size=24, sub_size=22, line_spacing_pts=16,
        bullets_width=Inches(6.6),
        extras=extras,
    )


# ---- Video 4 (77–84) ------------------------------------------------------

def slide_77_v4_title(prs):
    return make_video_title(prs, "Market Equilibrium", 4)


def slide_78_v4_outline(prs):
    return make_m1_outline(prs, 81, tag=TAG_V4, highlight_idx=2)


def slide_79_v4_mechanism(prs):
    def extras(slide):
        fig = SimpleFig(8.1, 6.40, 4.3, 4.05, 10, 10)
        _fig_axes(slide, fig, label_size=16)
        # D: y = 9 - x ; S: y = x ; equilibrium (4.5, 4.5)
        _fig_guide(slide, fig, (4.5, 4.5), color=GRAY)
        _add_arrow(slide, (fig.x(0), fig.y(6.5)), (fig.x(2.5), fig.y(6.5)),
                   color=GRAY, weight_pt=1.25, head=False, dash='dash')
        _add_arrow(slide, (fig.x(0), fig.y(2.5)), (fig.x(2.5), fig.y(2.5)),
                   color=GRAY, weight_pt=1.25, head=False, dash='dash')
        _fig_line(slide, fig, (1, 8), (8, 1), color=GOLD, weight_pt=2.75)
        _fig_curve_label(slide, fig, 8.05, 1.35, "D", color=NAVY)
        _fig_line(slide, fig, (1.5, 1.5), (8.5, 8.5), color=STEEL,
                  weight_pt=2.75)
        _fig_curve_label(slide, fig, 8.55, 8.9, "S", color=NAVY)
        # excess-supply band at P1 = 6.5 (D at 2.5, S at 6.5)
        _fig_line(slide, fig, (2.5, 6.5), (6.5, 6.5), color=BLUE_PED,
                  weight_pt=2.0)
        _add_text(slide, fig.x(2.9), fig.y(7.4), Inches(2.2), Inches(0.3),
                  "Excess supply", size=15, bold=True, color=BLUE_PED,
                  font="Calibri")
        # excess-demand band at P2 = 2.5 (S at 2.5, D at 6.5)
        _fig_line(slide, fig, (2.5, 2.5), (6.5, 2.5), color=GOLD,
                  weight_pt=2.0)
        _add_text(slide, fig.x(3.0), fig.y(2.2), Inches(2.4), Inches(0.3),
                  "Excess demand", size=15, bold=True, color=GOLD,
                  font="Calibri")
        _fig_ylab(slide, fig, 6.5, "P1", size=16)
        _fig_ylab(slide, fig, 4.5, "P0", size=16)
        _fig_ylab(slide, fig, 2.5, "P2", size=16)
        _fig_xlab(slide, fig, 4.5, "Q0", size=16)

    return content_slide(
        prs, 82, TAG_V4, "The Market Mechanism",
        [
            ([("Market equilibrium:", {'italic': True, 'color': RED}),
              (" where supply and demand intersect, at price ", {}),
              ("P0", {'italic': True}), (" and quantity ", {}),
              ("Q0", {'italic': True})], 0, {}),
            ([("At a higher price ", {}), ("P1", {'italic': True}),
              (", there is ", {}),
              ("excess supply", {'italic': True, 'color': BLUE_PED}),
              (", so the price falls", {})], 0, {}),
            ([("At a lower price ", {}), ("P2", {'italic': True}),
              (", there is ", {}),
              ("excess demand", {'italic': True, 'color': GOLD}),
              (", so the price is bid up", {})], 0, {}),
        ],
        size=24, sub_size=22, line_spacing_pts=18,
        bullets_width=Inches(6.6),
        extras=extras,
        notes=(
            "The point where they intersect is called market equilibrium, "
            "the price at which the quantity supplied equals the quantity "
            "demanded.\n"
            "To see why market equilibrium is stable, let’s go through "
            "what happens when the price deviates from the equilibrium "
            "condition:\n"
            "Suppose the price is higher than P_0; What happens at this "
            "price? Producers will start entering the market, wanting to "
            "sell at this high price, but they won’t all be able to find "
            "willing buyers. Then the quantity supplied will be greater "
            "than the quantity demanded. The excess quantity supplied is "
            "known as surplus. To eliminate the surplus, producers need to "
            "attract more buyers and to do so, they have to lower their "
            "prices. As price falls, quantity demanded increases, and the "
            "quantity supplied falls until the market reaches equilibrium.\n"
            "Suppose the price is lower than P_0; then the quantity "
            "demanded will exceed the quantity supplied and shortage will "
            "develop; To eliminate this shortage, buyers who cannot find "
            "the good available for sale will bid up the price and "
            "enterprising producers will be more than willing to raise "
            "their prices. As price rises, quantity demanded falls, "
            "quantity supplied rises, until we reach equilibrium.\n"
            "How does the market reach equilibrium in practice? Adam "
            "Smith’s invisible hand."),
    )


def slide_80_v4_terminology(prs):
    return content_slide(
        prs, 83, TAG_V4, "Summary of Terminology",
        [
            ("Equilibrium (or market clearing) price:", 0, {'bold': True}),
            ("Price that equates the quantity supplied to the quantity "
             "demanded", 1),
            ("Market mechanism:", 0, {'bold': True}),
            ("Tendency in a free market for price to change until the "
             "market clears", 1),
            ("Excess Supply:", 0, {'bold': True}),
            ("Situation in which the quantity supplied exceeds the "
             "quantity demanded", 1),
            ("Excess Demand (shortage):", 0, {'bold': True}),
            ("Situation in which the quantity demanded exceeds the "
             "quantity supplied", 1),
        ],
        size=24, sub_size=22, line_spacing_pts=12,
    )


def _v4_shift_chart(slide, *, d_shift=False, s_shift=False):
    """Shared chart for the three changes-in-equilibrium slides.
    Base: D: y = 9 - x, S: y = x (equilibrium 4.5, 4.5).
    D': y = 12 - x ; S': y = x - 2.5."""
    fig = SimpleFig(8.1, 6.40, 4.3, 4.05, 12, 12)
    _fig_axes(slide, fig, label_size=16)
    if d_shift and s_shift:
        q1, p1 = 7.25, 4.75
    elif d_shift:
        q1, p1 = 6.0, 6.0
    else:
        q1, p1 = 5.75, 3.25
    _fig_guide(slide, fig, (4.5, 4.5), color=GRAY)
    _fig_guide(slide, fig, (q1, p1), color=GRAY)
    _fig_line(slide, fig, (1, 8), (8, 1), color=GOLD, weight_pt=2.75)
    _fig_curve_label(slide, fig, 8.1, 1.4, "D", color=GOLD)
    _fig_line(slide, fig, (1.5, 1.5), (9.0, 9.0), color=STEEL,
              weight_pt=2.75)
    _fig_curve_label(slide, fig, 9.1, 9.4, "S", color=NAVY)
    if d_shift:
        _fig_line(slide, fig, (3.4, 8.6), (10.6, 1.4), color=GREEN_DK,
                  weight_pt=2.75, dash='dash')
        _fig_curve_label(slide, fig, 10.7, 1.9, "D’", color=GREEN_DK)
    if s_shift:
        _fig_line(slide, fig, (3.4, 0.9), (10.9, 8.4), color=BLUE_PED,
                  weight_pt=2.75, dash='dash')
        _fig_curve_label(slide, fig, 10.5, 8.7, "S’", color=BLUE_PED)
    # keep the two price labels legible when P1 sits close to P0
    p0_lab, p1_lab = 4.5, p1
    if abs(p1 - 4.5) < 0.7:
        p1_lab = 4.5 + 0.55 if p1 >= 4.5 else 4.5 - 0.55
        p0_lab = 4.5 - 0.55 if p1 >= 4.5 else 4.5 + 0.55
    _fig_ylab(slide, fig, p0_lab, "P0", size=16)
    _fig_ylab(slide, fig, p1_lab, "P1", size=16)
    _fig_xlab(slide, fig, 4.5, "Q0", size=16)
    _fig_xlab(slide, fig, q1, "Q1", size=16)
    return fig


def _v4_header(slide, runs):
    box = slide.shapes.add_textbox(int(Inches(7.3)), int(Inches(1.62)),
                                   int(Inches(5.3)), int(Inches(0.42)))
    tf = box.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    for text, opts in runs:
        r = p.add_run(); r.text = text
        r.font.name = "Calibri"; r.font.size = Pt(20)
        r.font.bold = True
        r.font.underline = opts.get('underline', False)
        r.font.color.rgb = NAVY
    return box


def slide_81_v4_shift_demand(prs):
    def extras(slide):
        _v4_header(slide, [("SHIFT IN DEMAND", {})])
        fig = _v4_shift_chart(slide, d_shift=True)
        # hand-tweaked 2026-08-23 (from logical 5.6,5.0 -> 7.0,6.0)
        _add_arrow(slide, (9208767, 4715447), (9667492, 4309110),
                   color=GREEN_DK, weight_pt=2.0, head=True)
        _add_text(slide, fig.x(7.2), fig.y(3.4), Inches(1.5), Inches(0.6),
                  "Shift in\ndemand", size=13, bold=True, color=GREEN_DK,
                  font="Calibri")

    return content_slide(
        prs, 84, TAG_V4, "Changes in Market Equilibrium",
        [
            ([("When the demand curve shifts to the right, the market "
               "clears at a higher price ", {}),
              ("P1", {'italic': True}),
              (" and a larger quantity ", {}),
              ("Q1", {'italic': True})], 0, {'bullet_style': 'none'}),
        ],
        size=24, sub_size=22, line_spacing_pts=0,
        bullets_width=Inches(6.4),
        extras=extras,
        notes=(
            "Shift out of demand implies that at any given price, "
            "consumers will want to buy more.\n"
            "If the price stayed the same, then a shortage would develop; "
            "price needs to rise in order to equilibrate the market.\n"
            "As price rises, quantity demanded falls and quantity supplied "
            "increases until we arrive at the new higher equilibrium price "
            "and quantity."),
    )


def slide_82_v4_shift_supply(prs):
    def extras(slide):
        _v4_header(slide, [("SHIFT IN SUPPLY", {})])
        fig = _v4_shift_chart(slide, s_shift=True)
        # arrow + label hand-tweaked 2026-08-23 (from logical
        # 6.2,6.6 -> 7.8,6.0 and label at 6.3, 7.7)
        _add_arrow(slide, (9634727, 3729609), (10044302, 4099941),
                   color=BLUE_PED, weight_pt=2.0, head=True)
        _add_text(slide, 9861803, 3551301, Inches(1.4), Inches(0.6),
                  "Shift in\nsupply", size=13, bold=True, color=BLUE_PED,
                  font="Calibri")

    return content_slide(
        prs, 85, TAG_V4, "Changes in Market Equilibrium",
        [
            ([("When the supply curve shifts to the right, the market "
               "clears at a lower price ", {}),
              ("P1", {'italic': True}),
              (" and a larger quantity ", {}),
              ("Q1", {'italic': True})], 0, {'bullet_style': 'none'}),
        ],
        size=24, sub_size=22, line_spacing_pts=0,
        bullets_width=Inches(6.4),
        extras=extras,
        notes=(
            "Shift out of supply implies that at any given price, "
            "producers are willing to sell a larger quantity.\n"
            "If the price stayed in its original position then there would "
            "be excess supply; price needs to fall to equilibrate the "
            "market.\n"
            "As the price falls, quantity demanded rises and quantity "
            "supplied declines.\n"
            "Price falls until the quantity supplied and quantity demanded "
            "are equal."),
    )


def slide_83_v4_shift_both(prs):
    def extras(slide):
        _v4_header(slide, [("SHIFT IN SUPPLY ", {}),
                           ("AND", {'underline': True}),
                           (" DEMAND", {})])
        _v4_shift_chart(slide, d_shift=True, s_shift=True)
        _add_ps_pointer(slide, left=MARGIN + Inches(0.2), top=Inches(6.35))

    return content_slide(
        prs, 86, TAG_V4, "Changes in Market Equilibrium",
        [
            ([("In this case, shifts in both curves lead to a slightly "
               "higher price ", {}),
              ("P1", {'italic': True}),
              (" and a much larger quantity ", {}),
              ("Q1", {'italic': True})], 0, {'bullet_style': 'none'}),
            ("Note: An even larger shift in S will lead to a lower price. "
             "But quantity unambiguously increases.", 0,
             {'bullet_style': 'none'}),
        ],
        size=24, sub_size=22, line_spacing_pts=18,
        bullets_width=Inches(6.4), bullets_height=Inches(4.2),
        extras=extras,
        notes=(
            "Two shifts simultaneously!\n"
            "When both curves shift, the prediction for either price or "
            "quantity will be ambiguous. It will depend on the magnitude "
            "of the shift in supply and demand. In the example on this "
            "slide, we observe that the prediction for price is ambiguous. "
            "It can either increase or decrease depending on how large the "
            "shifts in supply and demand are (above a small increase is "
            "depicted; as an exercise at home, see if you can shift the "
            "curves outward and generate a decrease in the equilibrium "
            "price). The prediction for quantity, however, is clear. It "
            "increases due to the shift out of both supply and demand."),
    )


def slide_84_shift_table(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_V4)
    _draw_action_title(
        slide, "Effect of Shifts in Demand and Supply Curves in Isolation")

    left, top = Inches(2.15), Inches(1.85)
    width, height = Inches(9.0), Inches(3.5)
    _add_graphicframe_shadow(slide, left - Inches(0.15), top - Inches(0.15),
                             width + Inches(0.3), height + Inches(0.3))
    gf = slide.shapes.add_table(6, 4, int(left), int(top),
                                int(width), int(height))
    tbl = gf.table
    tblPr = tbl._tbl.find(qn('a:tblPr'))
    if tblPr is not None:
        tblPr.set('firstRow', '0')
        tblPr.set('bandRow', '0')
        for child in list(tblPr):
            tblPr.remove(child)
    for i, w in enumerate([Inches(2.4), Inches(2.9), Inches(1.85),
                           Inches(1.85)]):
        tbl.columns[i].width = int(w)

    def _cell(r, c, text, *, fill, color, bold=True, size=18):
        cell = tbl.cell(r, c)
        cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.08)
        cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid(); cell.fill.fore_color.rgb = fill
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        if text:
            run = p.add_run(); run.text = text
            run.font.name = "Calibri"; run.font.size = Pt(size)
            run.font.bold = bold; run.font.color.rgb = color
        _set_cell_borders(cell)

    # header rows (with merges)
    _cell(0, 0, "Curve that shifts", fill=NAVY, color=WHITE)
    _cell(0, 1, "Direction of shift", fill=NAVY, color=WHITE)
    _cell(0, 2, "Impact on Equilibrium", fill=GOLD, color=NAVY)
    _cell(0, 3, "", fill=GOLD, color=NAVY)
    _cell(1, 0, "", fill=NAVY, color=WHITE)
    _cell(1, 1, "", fill=NAVY, color=WHITE)
    _cell(1, 2, "Price", fill=NAVY, color=WHITE)
    _cell(1, 3, "Quantity", fill=NAVY, color=WHITE)
    tbl.cell(0, 0).merge(tbl.cell(1, 0))
    tbl.cell(0, 1).merge(tbl.cell(1, 1))
    tbl.cell(0, 2).merge(tbl.cell(0, 3))

    body = [
        ("Demand Curve", "Out (increase in D)", "↑", "↑", CREAM),
        (None, "In (decrease in D)", "↓", "↓", WHITE),
        ("Supply Curve", "Out (increase in S)", "↓", "↑", CREAM),
        (None, "In (decrease in S)", "↑", "↓", WHITE),
    ]
    for i, (label, direction, p_arrow, q_arrow, fill) in enumerate(body):
        r = 2 + i
        _cell(r, 0, label or "", fill=RGBColor(0xEA, 0xEE, 0xF3),
              color=NAVY)
        _cell(r, 1, direction, fill=fill, color=NAVY, bold=False)
        _cell(r, 2, p_arrow, fill=fill, color=NAVY, size=22)
        _cell(r, 3, q_arrow, fill=fill, color=NAVY, size=22)
    tbl.cell(2, 0).merge(tbl.cell(3, 0))
    tbl.cell(4, 0).merge(tbl.cell(5, 0))

    # red-bordered "Important" rule box
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 int(Inches(1.4)), int(Inches(5.75)),
                                 int(Inches(10.5)), int(Inches(1.05)))
    try:
        box.adjustments[0] = 0.10
    except Exception:
        pass
    box.fill.solid(); box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = RED; box.line.width = Pt(1.5)
    box.shadow.inherit = False
    _add_drop_shadow(box)
    tb = slide.shapes.add_textbox(int(Inches(1.65)), int(Inches(5.82)),
                                  int(Inches(10.0)), int(Inches(0.9)))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r1 = p.add_run(); r1.text = "Important: "
    r1.font.name = "Calibri"; r1.font.size = Pt(16)
    r1.font.bold = True; r1.font.color.rgb = RED
    r2 = p.add_run()
    r2.text = ("As a general rule, when ")
    r3 = p.add_run(); r3.text = "both"
    r4 = p.add_run()
    # reworded by Nico 2026-08-23
    r4.text = (" curves shift the effect on equilibrium price and "
               "quantity depends on how much the curves shift.")
    for r in (r2, r3, r4):
        r.font.name = "Calibri"; r.font.size = Pt(16)
        r.font.italic = True; r.font.color.rgb = NAVY
    r3.font.underline = True

    _draw_footer(slide, FOOTER_TEXT, 87)
    # hidden by Nico 2026-08-23 (kept in the deck, skipped in the show)
    slide._element.set("show", "0")
    _set_notes(slide, (
        "Adopted from Melanie Wasserman's Module 1 deck: the four "
        "single-curve shifts and their price/quantity effects, plus the "
        "rule that simultaneous shifts pin down only one of the two "
        "directions with certainty. Good closing recap after the "
        "shift-in-supply-and-demand slide."))
    return slide


# ==========================================================================
#  2026-08-20 additions (approved by Nico): AC-heatwave solution (new #23,
#  from MW #51) and the copper mini-case (new #37–38, from MW #65–66).
#  NOTE: function names slide_NN_* above still carry the PRE-insert
#  numbering; the build() list below is the authority on display order.
# ==========================================================================

def slide_ac_solution(prs):
    """NEW #23 (from MW #51): solution to the AC-heatwave poll — the
    demand curve for AC shifts outward. Native D→D' chart."""
    def extras(slide):
        fig = SimpleFig(8.3, 6.35, 4.1, 3.9, 10, 10)
        _fig_axes(slide, fig, label_size=16)
        _fig_line(slide, fig, (1, 8), (8, 1), color=GOLD, weight_pt=2.75)
        _fig_curve_label(slide, fig, 8.05, 1.35, "D", color=GOLD)
        _fig_line(slide, fig, (3.4, 8.6), (10.0, 2.0), color=GREEN_DK,
                  weight_pt=2.75, dash='dash')
        # label + arrow hand-tweaked 2026-08-23 (from 9.65/2.6 and
        # 5.4->6.2; the D’ curve itself Nico left where it was)
        _fig_curve_label(slide, fig, 9.939, 2.7346, "D’", color=GREEN_DK)
        _add_arrow(slide, (fig.x(4.6), fig.y(4.8308)),
                   (fig.x(6.0), fig.y(5.6308)),
                   color=GREEN_DK, weight_pt=2.0, head=True)

    return content_slide(
        prs, 23, TAG_SD, "Shifts of the Demand Curve for AC",
        [
            ([("The heatwaves should cause the demand curve to ", {}),
              ("shift outward (to the right)",
               {'bold': True, 'color': GREEN_DK})], 0,
             {'bullet_style': 'none'}),
        ],
        size=26, sub_size=24, line_spacing_pts=0,
        bullets_width=Inches(6.6),
        extras=extras,
        notes=(
            "Adopted from Melanie Wasserman's Module 1 deck: closes the "
            "AC-heatwave poll. A heatwave is a non-price factor that "
            "raises the quantity demanded at every price, so the demand "
            "curve for air conditioners shifts outward, to the right."),
    )


def slide_copper_case(prs):
    """NEW #37 (from MW #65): the international copper market since 1880 —
    quantity up a hundredfold, real price flat."""
    def extras(slide):
        # MW's two-stage figure: quantity index first, price line painted
        # over it on the second bullet's click
        _add_media_image(slide, "mw_s65_rId3.png",
                         left=Inches(7.60), top=Inches(2.15),
                         width=Inches(4.95), rounded=False, shadow=True)
        _add_media_image(slide, "mw_s65_rId4.png",
                         left=Inches(7.60), top=Inches(2.15),
                         width=Inches(4.95), rounded=False, shadow=True)

    return content_slide(
        prs, 37, TAG_SD,
        "Mini-Case: The International Copper Market since 1880",
        [
            ("Annual consumption (in tons) increased hundredfold", 0),
            ("Price (inflation-adjusted) remained constant", 0),
            ("Use supply + demand framework to explain", 0),
        ],
        size=24, sub_size=22, line_spacing_pts=20,
        bullets_width=Inches(6.6), title_size=28,
        extras=extras,
        notes=(
            "Adopted from Melanie Wasserman's Module 1 deck. Copper "
            "consumption rose roughly a hundredfold since 1880 while the "
            "inflation-adjusted price stayed flat — the one shift "
            "combination the other mini-cases don't cover. Ask the class "
            "to explain it before showing the next slide."),
    )


def slide_copper_market(prs):
    """NEW #38 (from MW #66): both curves shift right — quantity soars,
    price flat."""
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_SD)
    _draw_action_title(slide, "The Market for Copper")
    fig = SimpleFig(4.6, 6.10, 5.3, 4.05, 12, 12)
    _sd_chart(
        slide, fig,
        curves=[
            ((1, 8), (8, 1), GOLD, None, "D0", (8.05, 1.35)),
            ((4, 9), (11.4, 1.6), GOLD, 'dash', "D1", (11.05, 2.1)),
            ((1.5, 1.5), (9, 9), STEEL, None, "S0", (9.05, 9.35)),
            ((5, 1), (11.5, 7.5), STEEL, 'dash', "S1", (11.2, 7.95)),
        ],
        points=[(4.5, 4.5), (8.5, 4.5)],
        xlabels=[(4.5, "Q0"), (8.5, "Q1")],
        ylabels=[(4.5, "P1 = P0")],
        arrows=[((3.9, 5.9), (5.3, 6.7), GRAY),
                ((6.4, 5.6), (7.8, 4.9), GRAY)],
    )
    # D1: y = 13 - x ; S1: y = x - 4 ; both meet at (8.5, 4.5) = (Q1, P0)
    _draw_footer(slide, FOOTER_TEXT, 38)
    _set_notes(slide, (
        "While it may looks puzzling at first, supply and demand shifts "
        "can explain it. Demand shifted right over time due to "
        "industrialization, population growth, and new uses "
        "(construction, electronics). Supply also shifted right because "
        "of advances in mining and transportation. These simultaneous "
        "shifts explain rising quantities but stable prices."))
    return slide


# ==========================================================================
# 2026-08-22 additions from "Module 1 - In Class with Solutions.pptx":
# 5 more PollEv stubs (Econ&Coffee pair + 3 results-view slides, each its
# own PollEv activity) and the 7-slide BACKUP section (displays 93-99)
# with jump links from slides 2/9/12/17 and back pills.
# ==========================================================================

def slide_poll_coffee_q(prs):
    return make_stub(prs, 7, TAG_LOG,
                     "Poll: Econ & Coffee weekend slot", STUB_POLL)


def slide_poll_coffee_results(prs):
    return make_stub(prs, 8, TAG_LOG,
                     "Poll results: Econ & Coffee weekend slot", STUB_POLL)


def slide_poll_ac_results(prs):
    return make_stub(prs, 25, TAG_SD,
                     "Poll results: heatwaves and AC demand", STUB_POLL)


def slide_poll_diamonds_results(prs):
    return make_stub(prs, 29, TAG_SD,
                     "Poll results: diamond demand", STUB_POLL)


def slide_poll_flip_results(prs):
    return make_stub(prs, 50, TAG_OPP,
                     "Poll results: flip-a-house profit", STUB_POLL)


def _draw_footer_nonum(slide):
    """Footer chrome without a page number (backup slides carry none)."""
    _add_rect(slide, 0, Inches(7.15), SLIDE_W, Inches(0.02), RULE)
    _add_rect(slide, MARGIN, Inches(7.135), GOLD_W, Inches(0.05), GOLD)
    _add_text(slide, MARGIN, Inches(7.20), Inches(11), Inches(0.32),
              FOOTER_TEXT, size=12, color=GRAY)


def _add_back_pill(slide, target_slide):
    """Navy '← Back' pill, lower-right (deck-standard position), jumping
    back to the slide that links here. Drawn AFTER the footer.

    2026-08-23: briefly carried an actionButtonBeginning chip; reverted to
    the original plain-text pill per Nico. The action-button family is used
    for the FORWARD jumps and the external-link markers only."""
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(11.72), Inches(6.60),
                                 Inches(1.55), Inches(0.46))
    try:
        shp.adjustments[0] = 0.5
    except Exception:
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = NAVY
    shp.line.fill.background()
    shp.shadow.inherit = False
    _add_drop_shadow(shp)
    tf = shp.text_frame
    tf.word_wrap = False
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "← Back"
    r.font.name = "Calibri"
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = WHITE
    shp.click_action.target_slide = target_slide
    return shp


def _add_click_overlay(slide, left, top, width, height, target_slide):
    """Invisible (100%-transparent-filled) rectangle carrying a
    jump-to-slide action. Run-level hyperlinks are unusable for styled
    text: this machine's PowerPoint renders hyperlinked runs UNDERLINED
    regardless of u="none" — verified 2026-08-22 against a native
    PowerPoint save, whose own no-underline hyperlink run also renders
    underlined. Transparent-fill (not noFill) keeps the interior
    hit-testable in the slideshow."""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 int(left), int(top),
                                 int(width), int(height))
    shp.fill.solid()
    shp.fill.fore_color.rgb = WHITE
    srgb = shp._element.spPr.find(qn('a:solidFill') + '/' + qn('a:srgbClr'))
    alpha = srgb.makeelement(qn('a:alpha'), {'val': '0'})
    srgb.append(alpha)
    shp.line.fill.background()
    shp.shadow.inherit = False
    shp.click_action.target_slide = target_slide
    return shp


def slide_93_backup_divider(prs):
    slide = _blank_slide(prs)
    _add_text(slide, 0, Inches(3.00), SLIDE_W, Inches(1.1), "BACKUP",
              size=54, bold=True, color=NAVY, font="Calibri",
              align=PP_ALIGN.CENTER)
    _add_rect(slide, int((SLIDE_W - Inches(4.0)) / 2), Inches(4.25),
              Inches(4.0), 54864, GOLD)
    _add_rect(slide, 0, Inches(7.15), SLIDE_W, Inches(0.02), RULE)
    _add_rect(slide, MARGIN, Inches(7.135), GOLD_W, Inches(0.05), GOLD)
    return slide


def slide_94_backup_leaders(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_BACKUP)
    _draw_action_title(slide, "Do National Leaders Matter?")
    # Econometrica header (Ottinger & Voigtländer 2025) + econimate still
    _add_media_image(slide, "ws63_image57.png",
                     left=Inches(0.65), top=Inches(1.55),
                     width=Inches(7.0), rounded=False, shadow=True)
    _add_media_image(slide, "ws63_image58.png",
                     left=Inches(7.75), top=Inches(3.55),
                     width=Inches(5.2), rounded=False, shadow=True)
    box = _add_hierarchical_bullets(
        slide, left=Inches(0.65), top=Inches(4.30),
        width=Inches(6.4), height=Inches(1.9),
        items=[
            # runs-list form: underline is a run-level option only
            ([("Economist Article",
               {'underline': True, 'bold': True})], 0, {}),
            ([("econimate Video (YouTube)",
               {'underline': True, 'bold': True})], 0, {}),
        ],
        size=22, line_spacing_pts=12)
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    _link_runs(slide, {
        "Economist Article":
            "https://www.economist.com/graphic-detail/2021/02/20/"
            "data-on-inbred-nobles-support-a-leader-driven-theory-of-history",
        "econimate Video (YouTube)":
            "https://www.youtube.com/watch?v=sOB5hmEXjdE",
    })
    # Document / Movie buttons after the two link labels. x from the
    # label ends in Calibri Bold 22 pt (3.244" / 4.567", text origin
    # 0.65" + 0.375" marL); y-centres 4.981" / 5.514" measured off the
    # outgoing gold glyphs in a 1500 px render.
    _add_ext_link_button(
        slide, "document", left=Inches(3.364), top=Inches(4.831),
        url="https://www.economist.com/graphic-detail/2021/02/20/"
            "data-on-inbred-nobles-support-a-leader-driven-theory-of-history")
    _add_ext_link_button(
        slide, "movie", left=Inches(4.687), top=Inches(5.364),
        url="https://www.youtube.com/watch?v=sOB5hmEXjdE")
    _draw_footer_nonum(slide)
    _add_back_pill(slide, prs.slides[1])          # -> display 2
    _set_notes(slide, (
        "Backup on my paper with Sebastian Ottinger, “History's "
        "Masters: The Effect of European Monarchs on State Performance” "
        "(Econometrica, January 2025). Using a millennium of European "
        "monarchs, we show that rulers' ability affected the performance "
        "of their states. The Economist covered the paper, and econimate "
        "made a short explainer video — both linked on the slide. "
        "The back button returns to the Introduction slide."))
    return slide


def slide_95_backup_happiness(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_BACKUP)
    _draw_action_title(
        slide, "Economics and a Common Critique: Does Money Buy Happiness?")
    box = _add_hierarchical_bullets(
        slide, left=MARGIN, top=Inches(1.50),
        width=Inches(6.5), height=Inches(0.6),
        items=[("The Easterlin Paradox (1974):", 0)],
        size=24, line_spacing_pts=0)
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    _add_media_image(slide, "ws64_image59.png",
                     left=Inches(3.02), top=Inches(2.40),
                     width=Inches(7.3), rounded=False, shadow=True)
    _draw_footer_nonum(slide)
    _set_notes(slide, (
        "First of two backup slides on money and happiness. Easterlin "
        "(1974) found that within a country at a point in time richer "
        "people report being happier, but average happiness did not seem "
        "to rise with GNP across countries — the “Easterlin "
        "paradox.” The next slide shows the Stevenson–Wolfers "
        "(2008) evidence that revisits this with better data."))
    return slide


def slide_96_backup_sw2008(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_BACKUP)
    _draw_action_title(slide, "Stevenson and Wolfers (2008)")
    _add_media_image(slide, "ws65_image60.jpeg",
                     left=Inches(3.90), top=Inches(1.50),
                     height=Inches(5.40), rounded=False, shadow=True)
    _draw_footer_nonum(slide)
    _add_back_pill(slide, prs.slides[11])         # -> display 12
    _set_notes(slide, (
        "Stevenson and Wolfers (2008) revisit the Easterlin paradox with "
        "richer cross-country data: life satisfaction rises with income "
        "both across and within countries, with no clear satiation point. "
        "The back button returns to the Homo Economicus slide."))
    return slide


def slide_97_backup_anderson(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_BACKUP)
    _draw_action_title(slide, "Related Work by Anderson Faculty")
    _add_media_image(slide, "ws66_goodlife_crop.png",
                     left=Inches(2.20), top=Inches(1.65),
                     width=Inches(8.9), rounded=False, shadow=True)
    box = _add_hierarchical_bullets(
        slide, left=Inches(2.20), top=Inches(5.35),
        width=Inches(9.4), height=Inches(1.2),
        items=[
            ("Note Dan Benjamin’s class “MGMT 298D-25 – "
             "Precision Healthcare”", 0),
            ("Applies concepts from (behavioral) economics", 1),
        ],
        size=22, sub_size=20, line_spacing_pts=6)
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    _draw_footer_nonum(slide)
    _add_back_pill(slide, prs.slides[11])         # -> display 12
    slide._element.set('show', '0')               # hidden in the source deck
    _set_notes(slide, (
        "Hidden backup: Anderson faculty work related to the "
        "money-and-happiness debate — Daniel Benjamin and co-authors' "
        "research on measuring well-being beyond GDP. His elective "
        "MGMT 298D-25 (Precision Healthcare) applies concepts from "
        "behavioral economics."))
    return slide


def slide_98_backup_portland(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_BACKUP)
    _draw_action_title(slide, "People Respond to Incentives")
    _add_media_image(slide, "ws67_image62.png",
                     left=Inches(0.90), top=Inches(1.62),
                     width=Inches(4.7), rounded=True, shadow=True)
    _add_text(slide, Inches(0.90), Inches(6.00), Inches(4.7), Inches(0.3),
              "Portland Street, Southampton, UK", size=12, italic=True,
              color=GRAY, font="Calibri", align=PP_ALIGN.CENTER)
    box = _add_hierarchical_bullets(
        slide, left=Inches(6.00), top=Inches(1.90),
        width=Inches(7.0), height=Inches(2.0),
        items=[
            ("Property tax:", 0, {'bold': True}),
            ([("Based on ", {}),
              ("the number of windows", {'bold': True})], 1, {}),
            ("England and Wales (1696 – 1851)", 1),
        ],
        size=24, sub_size=22, line_spacing_pts=8)
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    _add_media_image(slide, "ws67_image63.png",
                     left=Inches(7.30), top=Inches(4.05),
                     width=Inches(4.2), rounded=False, shadow=True)
    _draw_footer_nonum(slide)
    _add_back_pill(slide, prs.slides[8])          # -> display 9
    _set_notes(slide, (
        "This picture provides an excellent example on how to think like "
        "an economist. Do you notice anything strange? Some of the windows "
        "are missing, right? What could be going on? You may think that "
        "this is an error – maybe a terrible architect designed this "
        "building and skipped some windows. No, that’s not it. You "
        "may guess that this is a weird architectural style you had not "
        "heard about. Maybe the missing windows are a symbol for "
        "something, right? No, that’s not it either.\n\n"
        "So, what is it? Taxes, of course.\n\n"
        "In 1696, England introduced a property tax that was based on the "
        "number of windows. If your home has more than 10 windows, you "
        "have to pay more in taxes. The more windows, the higher the tax "
        "bill. This tax was intended to be a progressive tax: richer would "
        "pay more in taxes, because they tend to live in bigger houses, "
        "with more windows. This tax was also practical: a government "
        "employee could just walk to the front of your house, count the "
        "number of windows, and that is all of the information needed. No "
        "need to find out the number of bedrooms or bathrooms, or to "
        "figure out how much your home is worth.\n\n"
        "This tax had a terrible drawback though: it was easy to avoid. "
        "If you keep the number of windows below 10, you do not need to "
        "pay as much in taxes. This is the reason why many homes "
        "constructed in England during this period did not have as many "
        "windows. And homeowners would go as far as boarding up some of "
        "the existing windows to avoid the tax.\n\n"
        "Whoever came up with the idea for the windows tax did not think "
        "like an economist. By the end of this course, I hope you will "
        "not make the same mistake."))
    return slide


def slide_99_backup_prices(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_BACKUP)
    _draw_action_title(slide, "Can These Prices be Optimal?")
    _add_media_image(slide, "ws68_fares_top_crop.png",
                     left=Inches(2.17), top=Inches(1.42),
                     width=Inches(9.0), rounded=False, shadow=True)
    _add_media_image(slide, "ws68_fares_bot_crop.png",
                     left=Inches(2.17), top=Inches(4.50),
                     width=Inches(9.0), rounded=False, shadow=True)
    _draw_footer_nonum(slide)
    _add_back_pill(slide, prs.slides[16])         # -> display 17
    _set_notes(slide, (
        "Two Lufthansa bookings for the same November 2024 dates: on top, "
        "Los Angeles–Istanbul return connecting in Frankfurt "
        "($4,627 total); below, the nonstop Los Angeles–Frankfurt "
        "return — the Frankfurt legs are the same flights. Compare "
        "the totals in class and ask whether pricing like this can be "
        "optimal. The back button returns to the class agenda."))
    return slide


def wire_backup_links(prs):
    """Jump links from the main slides into the backup section. Runs after
    all slides exist. Slides 2 and 9 get a bare right-pointing action
    button at the end of the linked bullet line; slides 12 and 17 get a
    labelled jump pill with the same button in its left inset.

    2026-08-23: the invisible click overlays are gone. A transparent
    rectangle on top of the bullet box swallowed every click in that band,
    so the text box could not be selected or dragged in the editor. The
    action button is now the click target itself.

    The x positions below are the measured end of each linked line
    (PIL / Calibri: 9.72" at 24 pt, 11.20" at 28 pt, both from a text
    origin of 0.28" + 0.375" marL); the y positions are the line centres
    from the rendered layout. Re-check both after any font pass."""
    sl = lambda d: prs.slides[d - 1]
    _add_jump_button(sl(2), sl(96), left=Inches(9.80), top=Inches(4.555),
                     width=Inches(0.434), height=Inches(0.210))
    _add_jump_button(sl(9), sl(100), left=Inches(11.28), top=Inches(3.941),
                     width=Inches(0.490), height=Inches(0.238))
    # A backup link always takes the lower-RIGHT corner (Nico,
    # 2026-08-23); where a slide also carries a podcast / article link,
    # that one is placed wherever it fits best — on display 12 it is
    # centred in the space to the left.
    _add_jump_pill(sl(12), sl(97), left=Inches(8.99), top=Inches(6.58),
                   width=Inches(4.06), height=Inches(0.5),
                   label="Backup: Does money buy happiness?",
                   border=GOLD)
    # 2026-08-23: a backup link that is separate from the slide's own
    # text belongs in the lower RIGHT (Nico). 13.05" content edge - 4.10".
    _add_jump_pill(sl(17), sl(101), left=Inches(8.95), top=Inches(6.55),
                   width=Inches(4.10), height=Inches(0.5),
                   label="Backup: Can these prices be optimal?",
                   border=GOLD)


# ==========================================================================
# Build orchestration
# ==========================================================================

# --------------------------------------------------------------------------
# Economic symbols: italic letter + true subscript (2026-08-23, Nico)
# --------------------------------------------------------------------------
# "Whenever we have P0, Q0, etc, use subscript for the number ... any
# letters that symbolize price, quantity, etc need to be in italics, like
# in formulas."  Rather than touching ~40 label call sites, this runs as a
# deck-wide pass over every slide text run: any P / Q / D / S symbol
# followed by an index (a digit, or "Peak") is split into an italic base
# run and a subscript index run. The base letters are restricted to
# P, Q, D and S so ordinary text ("Module 1", "Fall 2026") is never touched.

SYMBOL_RE = re.compile(r"([PQDS])([\u2032\u2019']?)([0-9]|Peak)(?![A-Za-z0-9])")
SUBSCRIPT_BASELINE = "-25000"      # what PowerPoint writes for subscript


def _split_symbol_runs(para):
    """Rewrite the runs of one paragraph so every P/Q/D/S symbol carries an
    italic base and a subscript index."""
    changed = False
    for run in list(para.runs):
        text = run.text or ""
        if not SYMBOL_RE.search(text):
            continue
        pieces = []            # (text, is_subscript)
        pos = 0
        for m in SYMBOL_RE.finditer(text):
            if m.start() > pos:
                pieces.append((text[pos:m.start()], False))
            pieces.append((m.group(1) + m.group(2), False))
            pieces.append((m.group(3), True))
            pos = m.end()
        if pos < len(text):
            pieces.append((text[pos:], False))
        if len(pieces) < 2:
            continue
        r_el = run._r
        rPr = r_el.find(qn('a:rPr'))
        anchor = r_el
        for i, (txt, is_sub) in enumerate(pieces):
            if i == 0:
                run.text = txt
                if rPr is not None:
                    rPr.set('i', '1')
                continue
            new_r = copy.deepcopy(r_el)
            for t_el in new_r.findall(qn('a:t')):
                t_el.text = txt
            npr = new_r.find(qn('a:rPr'))
            if npr is not None:
                npr.set('i', '1')
                if is_sub:
                    npr.set('baseline', SUBSCRIPT_BASELINE)
                else:
                    npr.attrib.pop('baseline', None)
            anchor.addnext(new_r)
            anchor = new_r
        changed = True
    return changed


def apply_symbol_subscripts(prs):
    """Deck-wide: italic base + subscript index for every P/Q/D/S symbol."""
    n = 0
    for slide in prs.slides:
        for shp in slide.shapes:
            for tf in _iter_text_frames(shp):
                for para in tf.paragraphs:
                    if _split_symbol_runs(para):
                        n += 1
    return n


def _iter_text_frames(shp):
    if shp.has_text_frame:
        yield shp.text_frame
    if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shp.shapes:
            for tf in _iter_text_frames(child):
                yield tf
    if getattr(shp, "has_table", False) and shp.has_table:
        for row in shp.table.rows:
            for cell in row.cells:
                yield cell.text_frame


# --------------------------------------------------------------------------
# Speaker notes for the slides whose builders do not set their own
# (2026-08-23). Applied by apply_fill_notes() at the end of build(); a
# slide that already has notes is left alone, so source-ported notes and
# the PollEverywhere payload notes are never overwritten.
# --------------------------------------------------------------------------

FILL_NOTES = {
    1: "Welcome to Module 1. This module lays the foundation for the whole "
       "course: what a market is, how supply and demand set prices, and "
       "the three principles I want you to carry into every business "
       "decision. Let me start with a few words about myself and about how "
       "the class will run.",
    2: "A little about me. I trained as an environmental engineer before "
       "moving into economics, and my research asks why some countries are "
       "rich and others poor. That question is really about incentives and "
       "institutions, which is the same machinery we will use on business "
       "decisions in this course. My email is on the slide, and it is the "
       "right address for anything conceptual.",
    3: "Rafael is our TA and he is your first stop for anything "
       "procedural: problem sets, practice exercises and Achieve. "
       "Everything else lives on BruinLearn, including slides, handouts, "
       "the optional math review and the videos, and every class is "
       "recorded and livestreamed. Find those materials before the first "
       "problem set is due.",
    4: "The textbook is Goolsbee, Levitt and Syverson, fourth edition. "
       "Achieve is optional but useful if you want extra practice with "
       "worked feedback, and the direct link is on the slide. I will also "
       "post a weekly reading or podcast, and I am always glad to get "
       "suggestions from you.",
    5: "Five problem sets, worked in the study groups you have been "
       "assigned to, with the due dates on the class calendar. The grade "
       "weights are on the slide: midterm 35 percent, final 40 percent, "
       "problem sets 25 percent. Use AI freely while you are studying, but "
       "not in the exams.",
    6: "Send procedural questions to the TA and conceptual ones to me. The "
       "Econ and Coffee sessions are informal Zoom conversations where we "
       "start with concepts from class and then take problem-set "
       "questions. The TA also runs a recorded review session each week.",
    15: "You will get more out of this course by deciding where to invest "
        "your time, so set a realistic goal and hold to it. Focus on the "
        "intuition first. The math in this course is not hard, and I will "
        "always explain the economics behind it. Use your peers and use AI "
        "to push on your understanding, remembering that the exams are on "
        "your own.",
    17: "Here is the shape of the whole course. We are in part one, the "
        "basic principles and the economic way of thinking, and it feeds "
        "everything that follows: value and demand, then supply and cost, "
        "then markets, pricing and strategy. Keep this map in mind, "
        "because I will come back to it at each transition.",
    18: "These are the six topics of Module 1. The first three come from "
        "the videos: what a market is, how demand and supply work, and how "
        "they meet at an equilibrium. The last three are the principles we "
        "work through together in class. Problem Set 1 covers this "
        "material.",
    19: "You have watched the videos, so markets, demand and supply, and "
        "equilibrium are behind us. What we do together today is the three "
        "principles highlighted here, opportunity cost, sunk cost and "
        "marginal analysis, together with a set of mini-cases that put the "
        "supply-and-demand framework to work.",
    20: "A quick recall from Video 2. The practical test for whether two "
        "products sit in the same market is simple: if the price of the "
        "other product changes, does demand for yours move? That question "
        "is not academic. It decides antitrust cases, and we look at one "
        "next.",
    21: "In the 1990s ADM wanted to buy Clinton Corn Processing, and the "
        "Department of Justice objected because both had large shares of "
        "corn syrup. Everything turned on the market definition. DOJ said "
        "the market was corn syrup; ADM argued it was sweeteners more "
        "broadly, including sugar and artificial sweeteners, and used "
        "exactly the price test we just saw. ADM won in court, which is "
        "the lesson: defining the market is not a technicality, it is the "
        "case.",
    22: "Now try it yourselves. Is Netflix competing in online streaming, "
        "or in all film and television including theatres, or in all "
        "entertainment that competes for your evening? Each answer implies "
        "a different set of competitors and a different view of Netflix's "
        "market power, and there is no single right answer. That is the "
        "point.",
    23: "A Los Angeles heatwave, and a map of how many households have air "
        "conditioning. Before we draw anything, I want your instinct: what "
        "does a heatwave do to the market for air conditioners? Answer in "
        "the poll and we will see whether the class agrees.",
    31: "Tea is a clean example of both curves moving at once during "
        "Covid. People stuck at home drank more tea, and at the same time "
        "growing and shipping it got harder. Hold on to the two forces "
        "separately, because on the next slide we put them into one "
        "diagram.\n"
        "Note: Use this article starting in 2023: "
        "https://www.wsj.com/articles/lithium-prices-soar-turbocharged-by-"
        "electric-vehicle-demand-and-scant-supply-11639334956",
    32: "This is the same market drawn twice. Demand shifted out because "
        "people at home drank more tea, and supply shifted in because of "
        "bad weather, labor shortages and port closures. Both shifts push "
        "the price up, which is why the price increase here is "
        "unambiguous while the change in quantity is not.",
    33: "Think about the empty shelves you saw during the pandemic, or "
        "before a storm. Standard supply and demand says the price should "
        "rise until the market clears, so why do we see shortages instead? "
        "Part of the answer is that sellers worry about being seen as "
        "unfair, and that constraint is itself an economic force.",
    34: "Here is the news coverage that sets up the avocado case. Read the "
        "headline and hold one question in mind: what could make a price "
        "double and then halve inside a single year?",
    35: "The numbers are striking. From December 2016 to June 2017 the "
        "price of Mexican Hass avocados more than doubled, and from July "
        "to November 2017 it halved again. That pattern needs no special "
        "theory. Supply and demand handles it, as long as we are careful "
        "about which curve moves when.",
    36: "Three moves, in order. Demand shifted out with the guacamole and "
        "avocado-toast craze, taking us to D1. Then dry weather produced a "
        "meager crop and supply shifted left to S1, which is where the "
        "peak price comes from. Finally better pest control and the "
        "high-yielding half of the two-year cycle pushed supply back out "
        "to S2, so the price returned to where it started while the "
        "quantity ended much higher.",
    37: "Here is the recipe I want you to use every time. Name the shock, "
        "decide whether it hits demand or supply, then draw both curves "
        "before and after. Most mistakes in this course come from skipping "
        "the middle step.",
    38: "Apply the recipe to wheat. When the war began a major exporter "
        "was cut off, supply shifted left and the price jumped. When the "
        "Black Sea grain agreement reopened the export route, supply "
        "shifted back out and the price fell. Same market, two shocks, "
        "both on the supply side.",
    39: "Between September 2021 and September 2025, inflation-adjusted "
        "home prices in Los Angeles fell about 5 percent while the number "
        "of sales fell about 40 percent. A small price move next to a huge "
        "quantity move is a clue: this cannot be one curve shifting on its "
        "own. Work out which two curves moved, and in which direction, "
        "before we draw it.\n"
        "Source: https://www.redfin.com/city/11203/CA/Los-Angeles/"
        "housing-market",
    43: "We move now to the first of the three principles: economic costs "
        "include opportunity costs. This is the idea that costs you never "
        "write a check for still count.",
    44: "The opportunity cost of a choice is the value of the next best "
        "alternative you gave up. Economists count that alongside the "
        "money actually spent, which is why the full economic cost is the "
        "explicit cost plus the implicit one. Accountants see only the "
        "explicit part, and that difference is where a great many bad "
        "decisions live.",
    45: "A deliberately simple case. Given these alternatives, the "
        "opportunity cost of the one you pick is the value of the best one "
        "you did not. Notice that it depends on your own preferences, so "
        "two people facing identical options can face different "
        "opportunity costs.",
    46: "Three presents, each priced at 30 dollars, but worth different "
        "amounts to the person receiving them. The money cost is "
        "identical, so the whole decision turns on the value forgone. This "
        "is opportunity cost in its purest form.",
    47: "Now turn it on yourselves. Tuition is the explicit cost, but for "
        "most of you the larger number is the income and the career "
        "progress given up while studying. And here is a genuine question: "
        "does a remote or part-time option change that calculation?",
    48: "Here is the set-up. Buy the house for 500k, put in another 100k "
        "and a full year of your own time, and expect to sell for 700k. "
        "The alternative is a consulting job paying 150k that you cannot "
        "do at the same time. Before you answer the poll, be careful about "
        "which of those numbers belongs in the economic profit.",
    52: "There is one more opportunity cost hiding in that problem. The "
        "500k of capital could have been invested elsewhere and earned a "
        "return. Whether you count it depends on where the money came "
        "from, which is exactly the kind of judgment I want you to make "
        "explicit. There is a podcast on this, and we can pick it up at "
        "Econ and Coffee.",
    55: "The same pattern appears in US data estimated for 2022. The shape "
        "of the earnings path around the birth of a first child is "
        "strikingly similar across countries, which tells you this is not "
        "a quirk of one labor market. Every one of those forgone earnings "
        "is an opportunity cost.",
    56: "On to the second principle. Sunk costs are the ones you cannot "
        "get back, and the rule is to leave them out of the decision in "
        "front of you. It sounds obvious, and it is very hard to do.",
    57: "A sunk cost is one that has been paid and cannot be recovered. "
        "Because nothing you decide now can change it, it should not enter "
        "the decision at all. The hard part is psychological rather than "
        "analytical: walking away feels like admitting waste.",
    58: "Everyday versions are easy to find. Finishing dessert because the "
        "dinner was a fixed price, going to the gym because you paid for "
        "the year, skiing in bad weather because the ticket cannot be "
        "returned. Firms do the same thing with license fees and "
        "advertising already spent. We come back to this in Module 3.",
    59: "Concorde is the classic case. It was a genuine engineering "
        "triumph, New York to London in under three hours and a source of "
        "national pride for Air France and British Airways, and it lost "
        "enormous amounts of money. It never recovered its development "
        "cost and in the end did not even cover operating costs. Once that "
        "was clear the rational decision was to stop, and the pride is "
        "precisely what made stopping so hard.",
    60: "The rule is forward-looking. You cannot change the past, so "
        "optimize from here. Sunk costs usually show up because something "
        "unforeseen changed the right answer, and when that happens you "
        "re-optimize as though the money had never been spent.",
    61: "The third principle, and the one you will use most often: "
        "cost-benefit and marginal analysis. The question is never whether "
        "an activity is worth doing in the abstract, but whether the next "
        "unit of it is worth doing.",
    62: "The objective is to maximize net benefit, total benefit minus "
        "total cost. To find it you compare at the margin: the extra "
        "benefit from one more unit against the extra cost of that unit, "
        "opportunity costs included. Take the extra unit while marginal "
        "benefit exceeds marginal cost, and stop where the two are equal. "
        "MB equals MC is a rule you will meet in every module of this "
        "course.",
    63: "Let us do it with hours of exercise. The first hour carries a "
        "large net benefit, the second less, and by the fourth hour "
        "marginal benefit and marginal cost are equal, which is where you "
        "are indifferent. A fifth hour would cost more than it is worth, "
        "so four is the optimum. Notice that we never asked whether "
        "exercise is good, only whether the next hour is.",
    64: "The same logic drawn continuously, which is how you will see it "
        "for the rest of the course. Marginal benefit slopes down, "
        "marginal cost slopes up, and the optimum sits where they cross. "
        "Whenever this picture appears, it is the MB equals MC rule again.",
    65: "That is Module 1. Markets, supply and demand, and the equilibrium "
        "where they meet gave us the framework, and the three principles "
        "gave us the discipline. Economists count implicit costs where "
        "accountants do not, sunk costs get ignored, and decisions are "
        "made at the margin. Everything that follows builds on those "
        "ideas.",
    67: "This is the first of four short videos for Module 1. In this one "
        "I introduce the module and show how the pieces fit together.",
    68: "The same course map as in class. We are in part one, the basic "
        "principles and the economic way of thinking, and it is the "
        "foundation for value and demand, for supply and cost, and then "
        "for markets and pricing.",
    69: "Here are the six topics of Module 1 with a one-line description "
        "of each. The first three are covered in the videos, so you can "
        "watch them at your own pace. The last three are what we work "
        "through together in class.",
    70: "Video 2 is about markets: what a market is, and how you decide "
        "where its boundaries lie.",
    71: "Markets is the topic of this video. The question sounds simple "
        "and is not: which products, and which places, belong in the same "
        "market as yours?",
    72: "A company has to know three things about its market: who its "
        "customers are, who its competitors are, both actual and "
        "potential, and how far the market extends. The practical test for "
        "the product boundary is whether a price change elsewhere moves "
        "demand for your product. The geographic boundary varies "
        "enormously: a coffee shop in Venice competes with the next "
        "street, gasoline retail with the next few miles, gold with the "
        "whole world.",
    73: "Think about Netflix before class. Is the market online streaming, "
        "all films, or all entertainment? And did Covid change the answer, "
        "when theatres closed and streaming was the only option? Come with "
        "a view, because we will argue about it.",
    74: "Three kinds of actor, each with an objective. Consumers maximize "
        "utility from goods and services under a limited income. Workers "
        "choose jobs, trading off pay against leisure and flexibility. "
        "Firms hire those workers to produce for those consumers, and aim "
        "to maximize profit. Almost every model in this course is some "
        "combination of these three.",
    75: "Video 3 covers supply and demand, the workhorse framework of the "
        "whole course.",
    76: "This video is about demand and supply: how buyers and sellers "
        "each respond to price, and how to tell a movement along a curve "
        "from a shift of the curve itself.",
    77: "The question is how supply and demand together set the price and "
        "the quantity traded. The range of application is enormous: how "
        "growth in China affects the price of oil, of textiles, of "
        "electric cars. Once you can draw the two curves, you can reason "
        "about almost any market.",
    81: "This is the distinction students most often get wrong, so it is "
        "worth slowing down. A movement along D is a response to price "
        "with everything else held constant: the price rises from P1 to "
        "P2 and the quantity demanded falls from Q1 to Q2. A shift of D is "
        "caused by something other than price, such as higher income, and "
        "it moves the whole curve out to D prime, so at the same price P1 "
        "the quantity demanded is now Q3. Same curve, two very different "
        "events.",
    84: "The mirror image on the supply side. A movement along S is the "
        "response to price alone: the price rises from P1 to P2 and "
        "existing firms supply more, from Q1 to Q2. A shift of S comes "
        "from something else, such as better technology, cheaper inputs, "
        "or new firms entering, and it moves the whole curve out to S "
        "prime, so at the same price P2 the quantity supplied is Q3.",
    85: "Video 4 brings the two sides together: market equilibrium, and "
        "what happens when either curve moves.",
    86: "Equilibrium is the topic here: where demand meets supply, what "
        "makes the market get there, and what happens when one of the "
        "curves shifts.",
    88: "Four terms you will need. The equilibrium, or market-clearing, "
        "price is the one at which quantity supplied equals quantity "
        "demanded. The market mechanism is the tendency for price to move "
        "until that happens. Above that price there is excess supply, "
        "below it excess demand, which is a shortage. Keep the words "
        "straight and the diagrams become much easier.",
    93: "Everything from here is backup material: slides I keep for "
        "questions that come up in class, and the targets of the links "
        "earlier in the deck. Nothing after this point is part of the main "
        "flow.",
}

# these two carry a source link in their existing notes, which the new
# text above reproduces -- so they are allowed to overwrite
FILL_NOTES_OVERWRITE = {31, 39}


def _m1_disp_shift(k):
    """2026-08-23: the two Tapestry slides were inserted at displays
    73-74, so every display from 73 on moved down by two."""
    return k + 2 if k >= 73 else k


def apply_fill_notes(prs):
    """Fill in notes for slides whose builders set none. Existing notes
    are never overwritten (except the two above), which keeps the
    source-ported notes and the PollEverywhere payload notes intact."""
    n = 0
    for key, text in FILL_NOTES.items():
        disp = _m1_disp_shift(key)
        slide = prs.slides[disp - 1]
        if key not in FILL_NOTES_OVERWRITE:
            if slide.has_notes_slide:
                existing = slide.notes_slide.notes_text_frame.text or ""
                if existing.strip():
                    continue
        _set_notes(slide, text)
        n += 1
    return n


def build(out_path=None):
    prs = Presentation()
    prs.slide_width = int(SLIDE_W)
    prs.slide_height = int(SLIDE_H)

    slide_01_title(prs)                                            #  1
    slide_02_introduction(prs)                                     #  2
    slide_03_logistics1(prs)                                       #  3
    slide_04_logistics2(prs)                                       #  4
    slide_05_logistics3(prs)                                       #  5
    slide_06_office_hours(prs)                                     #  6
    slide_poll_coffee_q(prs)                                       #  7 NEW poll
    slide_poll_coffee_results(prs)                                 #  8 NEW poll
    slide_07_why_econ(prs)                                         #  9
    slide_08_models(prs)                                           # 10
    slide_09_find_model(prs)                                       # 11
    slide_10_homo_economicus(prs)                                  # 12
    slide_11_hedgehogs(prs)                                        # 13
    slide_12_making_most_1(prs)                                    # 14
    slide_13_making_most_2(prs)                                    # 15
    slide_14_teaching_philosophy(prs)                              # 16
    make_roadmap(prs, 17)                                          # 17
    slide_16_outline(prs)                                          # 18
    slide_17_outline_now(prs)                                      # 19
    slide_18_recall_market_def(prs)                                # 20
    slide_19_adm(prs)                                              # 21
    slide_20_netflix(prs)                                          # 22
    slide_21_heatwaves(prs)                                        # 23
    slide_22_poll_ac(prs)                                          # 24
    slide_poll_ac_results(prs)                                     # 25 NEW poll
    slide_ac_solution(prs)                                         # 26
    slide_23_swiftonomics(prs)                                     # 27
    slide_24_poll_diamonds(prs)                                    # 28
    slide_poll_diamonds_results(prs)                               # 29 NEW poll
    slide_25_swift_solution(prs)                                   # 30
    slide_26_tea(prs)                                              # 31
    slide_27_tea_market(prs)                                       # 32
    slide_28_disasters(prs)                                        # 33
    slide_29_avocado_clip(prs)                                     # 34
    slide_30_avocado_bullets(prs)                                  # 35
    slide_31_avocado_market(prs)                                   # 36
    slide_32_steps(prs)                                            # 37
    slide_33_wheat(prs)                                            # 38
    slide_34_la_case(prs)                                          # 39
    slide_35_la_market(prs)                                        # 40
    slide_copper_case(prs)                                         # 41
    slide_copper_market(prs)                                       # 42
    slide_36_outline_opp(prs)                                      # 43
    slide_37_opp_costs(prs)                                        # 44
    slide_38_fruit_table(prs)                                      # 45
    slide_39_present(prs)                                          # 46
    slide_40_mba_cost(prs)                                         # 47
    slide_41_flip_house(prs)                                       # 48
    slide_42_poll_flip(prs)                                        # 49
    slide_poll_flip_results(prs)                                   # 50 NEW poll
    slide_43_flip_solution(prs)                                    # 51
    slide_44_another_opp(prs)                                      # 52
    slide_45_child_cost(prs)                                       # 53
    slide_46_child_penalty(prs)                                    # 54
    slide_47_us_2022(prs)                                          # 55
    slide_48_outline_sunk(prs)                                     # 56
    slide_49_sunk_costs(prs)                                       # 57
    slide_50_sunk_examples(prs)                                    # 58
    slide_51_concorde(prs)                                         # 59
    slide_52_sunk_takeaway(prs)                                    # 60
    slide_53_outline_cba(prs)                                      # 61
    slide_54_cba(prs)                                              # 62
    slide_55_exercise(prs)                                         # 63
    slide_56_continuous(prs)                                       # 64
    slide_57_summary(prs)                                          # 65
    slide_58_next_steps(prs)                                       # 66
    slide_59_v1_title(prs)                                         # 67
    slide_60_v1_roadmap(prs)                                       # 68
    slide_61_v1_outline(prs)                                       # 69
    slide_62_v2_title(prs)                                         # 70
    slide_63_v2_outline(prs)                                       # 71
    slide_64_v2_market_def(prs)                                    # 72
    slide_tapestry_case(prs)                                       # 73 NEW
    slide_tapestry_evidence(prs)                                   # 74 NEW
    slide_65_v2_netflix(prs)                                       # 73
    slide_66_v2_actors(prs)                                        # 74
    slide_67_v3_title(prs)                                         # 75
    slide_68_v3_outline(prs)                                       # 76
    slide_69_v3_ds_analysis(prs)                                   # 77
    slide_70_v3_demand_def(prs)                                    # 78
    slide_71_v3_ceteris_paribus(prs)                               # 79
    slide_72_v3_demand_curve(prs)                                  # 80
    slide_73_v3_move_vs_shift_d(prs)                               # 81
    slide_74_v3_ai_chips(prs)                                      # 82
    slide_75_v3_supply_curve(prs)                                  # 83
    slide_76_v3_move_vs_shift_s(prs)                               # 84
    slide_77_v4_title(prs)                                         # 85
    slide_78_v4_outline(prs)                                       # 86
    slide_79_v4_mechanism(prs)                                     # 87
    slide_80_v4_terminology(prs)                                   # 88
    slide_81_v4_shift_demand(prs)                                  # 89
    slide_82_v4_shift_supply(prs)                                  # 90
    slide_83_v4_shift_both(prs)                                    # 91
    slide_84_shift_table(prs)                                      # 92
    slide_93_backup_divider(prs)                                   # 93 BACKUP
    slide_94_backup_leaders(prs)                                   # 94 BACKUP
    slide_95_backup_happiness(prs)                                 # 95 BACKUP
    slide_96_backup_sw2008(prs)                                    # 96 BACKUP
    slide_97_backup_anderson(prs)                                  # 97 BACKUP (hidden)
    slide_98_backup_portland(prs)                                  # 98 BACKUP
    slide_99_backup_prices(prs)                                    # 99 BACKUP

    wire_backup_links(prs)
    n_notes = apply_fill_notes(prs)
    n_sym = apply_symbol_subscripts(prs)

    out = Path(out_path) if out_path else OUT_DIR / "Module 1 - Revised.pptx"
    prs.save(str(out))
    print(f"saved {out} — {len(prs.slides._sldIdLst)} slides "
          f"({n_sym} paragraph(s) with subscripted symbols, "
          f"{n_notes} slide(s) given fill-in notes)")
    return out


if __name__ == "__main__":
    import sys as _sys
    build(_sys.argv[1] if len(_sys.argv) > 1 else None)
