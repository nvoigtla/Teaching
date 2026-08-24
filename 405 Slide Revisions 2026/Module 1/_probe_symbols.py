# -*- coding: utf-8 -*-
"""Probe: the whole PowerPoint action-button family in the deck palette
(gold face, navy glyph + hairline border), at the size used on slide 2 —
so the fourth type, the one marking "this links to another slide", can be
picked by eye. There is no `actionButtonSlide` preset; these 12 are all
PowerPoint offers.
"""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from _build_template_samples import (
    GOLD, GRAY, NAVY, RULE, SLIDE_H, SLIDE_W, WHITE,
    _add_rect, _add_text, _blank_slide,
)
from _build_Module1 import _add_drop_shadow

OUT = "_probe_symbols.pptx"

FAMILY = [
    ("END", MSO_SHAPE.ACTION_BUTTON_END, "in use: slides 2, 9, 12, 17"),
    ("BEGINNING", MSO_SHAPE.ACTION_BUTTON_BEGINNING, "in use: Back pills"),
    ("SOUND", MSO_SHAPE.ACTION_BUTTON_SOUND, "in use: podcast"),
    ("DOCUMENT", MSO_SHAPE.ACTION_BUTTON_DOCUMENT, "in use: article"),
    ("MOVIE", MSO_SHAPE.ACTION_BUTTON_MOVIE, "in use: video"),
    ("INFORMATION", MSO_SHAPE.ACTION_BUTTON_INFORMATION, "candidate"),
    ("FORWARD_OR_NEXT", MSO_SHAPE.ACTION_BUTTON_FORWARD_OR_NEXT, "candidate"),
    ("BACK_OR_PREVIOUS", MSO_SHAPE.ACTION_BUTTON_BACK_OR_PREVIOUS, ""),
    ("RETURN", MSO_SHAPE.ACTION_BUTTON_RETURN, ""),
    ("HOME", MSO_SHAPE.ACTION_BUTTON_HOME, ""),
    ("HELP", MSO_SHAPE.ACTION_BUTTON_HELP, ""),
    ("CUSTOM", MSO_SHAPE.ACTION_BUTTON_CUSTOM, "blank face"),
]


def btn(slide, shape, left, top, w=Inches(0.62), h=Inches(0.30)):
    shp = slide.shapes.add_shape(shape, int(left), int(top), int(w), int(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = GOLD
    shp.line.color.rgb = NAVY
    shp.line.width = Pt(1.25)
    shp.shadow.inherit = False
    _add_drop_shadow(shp)
    return shp


def line(slide, y, text, shape, size=24):
    b = slide.shapes.add_textbox(int(Inches(0.65)), int(y),
                                 int(Inches(9.2)), int(Inches(0.42)))
    tf = b.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = "Calibri"
    r.font.size = Pt(size)
    r.font.color.rgb = NAVY
    if shape is not None:
        btn(slide, shape, Inches(9.80), y + Inches(0.06))


def main():
    prs = Presentation()
    prs.slide_width = int(SLIDE_W)
    prs.slide_height = int(SLIDE_H)

    # ---- sheet 1: the whole family ------------------------------------
    s = _blank_slide(prs)
    _add_rect(s, 0, 0, SLIDE_W, Inches(0.42), NAVY)
    _add_text(s, Inches(0.28), 0, Inches(12), Inches(0.42),
              "The 12 PowerPoint action buttons, in the deck palette "
              "(there is no actionButtonSlide)",
              size=16, bold=True, color=WHITE)
    _add_rect(s, Inches(0.28), Inches(0.90), Inches(12.78), Inches(0.02),
              RULE)
    for i, (name, shape, note) in enumerate(FAMILY):
        col = i // 6
        row = i % 6
        x = Inches(0.55) + col * Inches(6.5)
        y = Inches(1.25) + row * Inches(0.92)
        btn(s, shape, x, y, w=Inches(0.78), h=Inches(0.38))
        _add_text(s, x + Inches(1.05), y - Inches(0.03), Inches(3.1),
                  Inches(0.32), name, size=16, bold=True, color=NAVY)
        if note:
            _add_text(s, x + Inches(1.05), y + Inches(0.20), Inches(4.2),
                      Inches(0.28), note, size=12, color=GRAY)
    _add_rect(s, 0, Inches(7.15), SLIDE_W, Inches(0.02), RULE)

    # ---- sheet 2: the three candidates on slide 2's actual line --------
    s2 = _blank_slide(prs)
    _add_rect(s2, 0, 0, SLIDE_W, Inches(0.42), NAVY)
    _add_text(s2, Inches(0.28), 0, Inches(12), Inches(0.42),
              "Candidates for the slide-jump marker, on slide 2's line "
              "(24 pt)", size=16, bold=True, color=WHITE)
    _add_rect(s2, Inches(0.28), Inches(0.90), Inches(12.78), Inches(0.02),
              RULE)
    txt = "My research: Why are some countries so rich and others so poor?"
    rows = [("END  —  currently in the deck", MSO_SHAPE.ACTION_BUTTON_END),
            ("INFORMATION  —  'there is more on a backup slide'",
             MSO_SHAPE.ACTION_BUTTON_INFORMATION),
            ("FORWARD_OR_NEXT  —  plain triangle, boxed",
             MSO_SHAPE.ACTION_BUTTON_FORWARD_OR_NEXT),
            ("DOCUMENT  —  reads as 'a page', already used for articles",
             MSO_SHAPE.ACTION_BUTTON_DOCUMENT)]
    for i, (label, shape) in enumerate(rows):
        y = Inches(1.35) + i * Inches(1.40)
        _add_text(s2, Inches(0.65), y, Inches(11.0), Inches(0.30),
                  label, size=14, bold=True, color=GRAY)
        line(s2, y + Inches(0.38), txt, shape)
    _add_rect(s2, 0, Inches(7.15), SLIDE_W, Inches(0.02), RULE)

    prs.save(OUT)
    print("saved", OUT)


if __name__ == "__main__":
    main()
