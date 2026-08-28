# ==========================================================================
#  _build_Module2InClass.py — phase-1 scaffold for
#  "Module 2 - In Class Revised.pptx"
#
#  Demand Analysis, in-class deck (76 slides per the approved outline in
#  "Module 2 - In Class Revised - outline.md"; the 16 PollEverywhere
#  slides + the pizza Excel-embed slide are positional stubs until the
#  _splice_media.py pass).
#
#  Helper layer copied VERBATIM from Module 7/_build_Module7.py
#  (2026-08-14), which itself carries the Module 3 helper layer — the
#  proven chrome, box, bullet, OMML, chart, and table primitives.
#  M2-specific code starts at the "MODULE 2" banner below.
# ==========================================================================

import copy
import math
import re
import shutil
import zipfile
from pathlib import Path

from lxml import etree as ET
from pptx import Presentation
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Cm, Inches, Pt

# Reuse all primitives from the template script (single source of truth).
from _build_template_samples import (
    FADED,
    FOOTER_TEXT,
    GOLD,
    GOLD_W,
    GRAY,
    MARGIN,
    MODULE_AGENDA,
    NAVY,
    RULE,
    RULE_W,
    SLIDE_H,
    SLIDE_W,
    WHITE,
    _add_bulleted_list,
    _add_rect,
    _add_text,
    _blank_slide,
    _draw_action_title,
    _set_bullet_char,
)

OUT_DIR = Path(__file__).parent


# --------------------------------------------------------------------------
# Title-case top bar – replaces the all-caps default from the template script.
# The user prefers academic-paper title case (each major word capitalized,
# small connectors like "and", "of", "for", "the" stay lowercase).
# --------------------------------------------------------------------------

def _draw_top_bar_tc(slide, section_tag):
    """Navy top bar with title-cased section tag (no uppercase forcing)."""
    bar_h = Inches(0.42)
    _add_rect(slide, 0, 0, SLIDE_W, bar_h, NAVY)
    _add_text(slide, MARGIN, 0, Inches(12), bar_h,
              section_tag, size=16, bold=True,
              color=WHITE, font="Calibri",
              anchor=MSO_ANCHOR.MIDDLE)


def _draw_footer(slide, footer_text, page_num):
    """Footer rule + gold accent + footer text/page number, with larger type
    than the template default so handout printouts remain legible."""
    _add_rect(slide, 0, Inches(7.15), SLIDE_W, Inches(0.02), RULE)
    _add_rect(slide, MARGIN, Inches(7.135), GOLD_W, Inches(0.05), GOLD)
    _add_text(slide, MARGIN, Inches(7.20), Inches(11), Inches(0.32),
              footer_text, size=12, color=GRAY)
    _add_text(slide, Inches(12.5), Inches(7.20), Inches(0.6), Inches(0.32),
              str(page_num), size=12, color=GRAY, align=PP_ALIGN.RIGHT)


# --------------------------------------------------------------------------
# Speaker-notes helper
# --------------------------------------------------------------------------

def _set_notes(slide, text):
    """Replace the slide's speaker notes with *text*."""
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.clear()
    notes_tf.text = text


# --------------------------------------------------------------------------
# Reusable shape primitives for diagrams
# --------------------------------------------------------------------------

def _add_filled_box(slide, left, top, width, height, label, *,
                    fill=NAVY, text_color=WHITE, line=None,
                    size=18, bold=True, font="Calibri"):
    """Filled rectangle with centered text."""
    left, top, width, height = int(left), int(top), int(width), int(height)
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height,
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = text_color
    return shp


def _add_oval_outline(slide, left, top, width, height, *,
                      color=GOLD, weight_pt=2.25, shadow=True):
    """Unfilled oval used to ring a symbol inside a formula (2026-08-25,
    Nico's circled percentage sign on slide 28, copied from CT)."""
    shp = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, int(left), int(top), int(width), int(height))
    shp.fill.background()
    shp.line.color.rgb = color
    shp.line.width = Pt(weight_pt)
    shp.shadow.inherit = False
    if shadow:
        _add_drop_shadow(shp)
    return shp


def _add_rot_brace(slide, left, top, width, height, rot_deg, *,
                   color=GOLD, weight_pt=2.0):
    """A right-brace rotated to lie along a sloped line (CT's device on
    her In-Class slide 39: one brace per stretch of the demand curve,
    open side onto the line, nub pointing at the label)."""
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_BRACE, int(left), int(top), int(width),
        int(height))
    shp.fill.background()
    shp.line.color.rgb = color
    shp.line.width = Pt(weight_pt)
    shp.shadow.inherit = False
    shp.rotation = rot_deg
    return shp


def _add_oval_filled(slide, left, top, width, height, *,
                     fill=NAVY, line=None, weight_pt=1.0):
    """A filled dot / marker placed at absolute coordinates (the
    figure-relative version is ``_fig_point``)."""
    shp = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, int(left), int(top), int(width), int(height))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(weight_pt)
    shp.shadow.inherit = False
    return shp


def _add_runs_text(slide, left, top, width, height, runs, *,
                   size=18, bold=False, italic=False, color=NAVY,
                   font="Calibri", align=PP_ALIGN.LEFT):
    """One-line text box whose words carry different formatting.

    ``runs`` is a list of ``(text, opts)`` with ``opts`` keys
    ``size`` / ``bold`` / ``italic`` / ``color``; anything omitted falls
    back to the box-level default."""
    box = slide.shapes.add_textbox(int(left), int(top), int(width),
                                   int(height))
    tf = box.text_frame
    tf.word_wrap = False
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    for text, opts in runs:
        r = p.add_run()
        r.text = text
        r.font.name = font
        r.font.size = Pt(opts.get('size', size))
        r.font.bold = opts.get('bold', bold)
        r.font.italic = opts.get('italic', italic)
        r.font.color.rgb = opts.get('color', color)
    return box


def _add_outlined_box(slide, left, top, width, height, label, *,
                      line=NAVY, text_color=NAVY, fill=WHITE,
                      size=18, bold=True, line_w=1.25, font="Calibri",
                      rounded=False, shadow=False, corner_pct=0.06,
                      sub_label=None, sub_size=None):
    """Outlined rectangle (white fill) with centered text.

    ``rounded=True`` switches the base shape to a rounded rectangle
    (corner-adjust = ``corner_pct``); ``shadow=True`` adds a soft drop
    shadow.  Both default off so existing flat call-sites are
    unaffected.
    """
    left, top, width, height = int(left), int(top), int(width), int(height)
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(
        shape_type, left, top, width, height,
    )
    if rounded:
        try: shp.adjustments[0] = corner_pct
        except Exception: pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    if shadow:
        _add_drop_shadow(shp)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = text_color
    if sub_label:
        # a second, smaller line under the label - used by the practice
        # video reference boxes (2026-08-25)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = sub_label
        r2.font.name = font
        r2.font.size = Pt(sub_size if sub_size else max(11, size - 3))
        r2.font.bold = False
        r2.font.color.rgb = text_color
    return shp


def _add_convention_box(slide, left, top, width, height, *,
                          prefix=None, body=None, runs=None,
                          fill_rgb=None, border=None, line_w=1.0,
                          corner_pct=0.12, size=15, align=PP_ALIGN.LEFT,
                          font="Calibri", pad_h=None, pad_v=None,
                          line_spacing_pct=None, fill_alpha_pct=None):
    """Cream-fill / navy-border rounded-rect explanation callout.

    The "Convention" textbox pattern from slide 14 generalised — use it
    anywhere a slide needs a compact, visually-distinct box for a
    short conceptual explanation or notational convention.  Sits well
    below a table, beside a hero formula, or as a slide-wide footer.

    Two ways to populate the text:
      • ``prefix`` (bold) + ``body`` (regular) — simplest path; matches
        slide 14's "Convention:  <text>" pattern.
      • ``runs`` — a list of ``(text, {"bold": .., "italic": .., ...})``
        tuples for finer-grained styling (multi-line, mixed formatting).

    Style defaults follow the course-layer CLAUDE.md "Convention callout
    box" spec — cream fill, thin primary-color border, slight rounding,
    primary-color text.  Override ``fill_rgb`` / ``border`` only when you
    need a different accent.

    ``fill_alpha_pct`` washes the fill out to that opacity (2026-08-26,
    Nico): the WARNING variant of this box is a transparent dark-red
    tint with a dark-red border, used where a line is a caution rather
    than a convention.
    """
    fill = fill_rgb if fill_rgb is not None else RGBColor(0xFD, 0xF6, 0xE6)
    border = border if border is not None else NAVY

    left, top, width, height = int(left), int(top), int(width), int(height)
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if fill_alpha_pct is not None:
        srgb = shp._element.spPr.find(
            qn('a:solidFill')).find(qn('a:srgbClr'))
        alpha = ET.SubElement(srgb, qn('a:alpha'))
        alpha.set('val', str(int(fill_alpha_pct * 1000)))
    shp.line.color.rgb = border
    shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    try:
        shp.adjustments[0] = corner_pct
    except Exception:
        pass
    # 2026-08-24 (Nico): the cream callouts carry the deck's soft shade
    # like every other filled box, so they read as lifted cards
    _add_drop_shadow(shp)

    # Inset text box so the rounded corners breathe — matches slide 14.
    pad_h = Inches(0.20) if pad_h is None else pad_h
    pad_v = Inches(0.12) if pad_v is None else pad_v
    tb = slide.shapes.add_textbox(
        left + int(pad_h), top + int(pad_v),
        width - 2 * int(pad_h), height - 2 * int(pad_v),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0); tf.margin_bottom = Inches(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    def _style_run(r, opts):
        r.font.name = opts.get('font', font)
        r.font.size = Pt(opts.get('size', size))
        r.font.bold = opts.get('bold', False)
        r.font.italic = opts.get('italic', False)
        r.font.underline = opts.get('underline', False)
        r.font.color.rgb = opts.get('color', NAVY)

    if runs is not None:
        first = True
        for entry in runs:
            text, opts = entry if isinstance(entry, tuple) else (entry, {})
            if opts.get('newline') and not first:
                p = tf.add_paragraph()
            elif first:
                p = tf.paragraphs[0]
            else:
                # Same paragraph — append run to the most-recent paragraph.
                p = tf.paragraphs[-1]
            p.alignment = align
            r = p.add_run()
            r.text = text
            _style_run(r, opts)
            first = False
    else:
        p = tf.paragraphs[0]
        p.alignment = align
        if prefix:
            r1 = p.add_run(); r1.text = prefix
            _style_run(r1, {'bold': True, 'color': NAVY, 'size': size})
        if body:
            r2 = p.add_run(); r2.text = body
            _style_run(r2, {'color': NAVY, 'size': size})

    if line_spacing_pct is not None:
        for p_obj in tf.paragraphs:
            pPr = p_obj._p.get_or_add_pPr()
            for old in pPr.findall(qn('a:lnSpc')):
                pPr.remove(old)
            lnSpc = ET.Element(qn('a:lnSpc'))
            spcPct = ET.SubElement(lnSpc, qn('a:spcPct'))
            spcPct.set('val', str(int(line_spacing_pct * 1000)))
            pPr.insert(0, lnSpc)
    return shp


def _add_rounded_filled_box(slide, left, top, width, height, label, *,
                             fill=NAVY, text_color=WHITE, line=None,
                             size=18, bold=True, italic=False,
                             font="Calibri",
                             corner_pct=0.06, shadow=True, line_w=0.75):
    """Rounded-corner filled rectangle with centered text and soft drop shadow.

    Mirrors :func:`_add_filled_box` but renders ``MSO_SHAPE.ROUNDED_RECTANGLE``
    with the corner-adjust set to ``corner_pct`` (6 % per course CLAUDE.md
    "slight rounding") and a soft drop shadow via :func:`_add_drop_shadow`.
    """
    left, top, width, height = int(left), int(top), int(width), int(height)
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
    )
    try:
        shp.adjustments[0] = corner_pct
    except Exception:
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    if shadow:
        _add_drop_shadow(shp)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = text_color
    return shp


def _add_arrow(slide, start_xy, end_xy, *, color=NAVY, weight_pt=1.5,
               head=True, dash=None, head_size='med'):
    """Draw a line/arrow from start to end (in EMU/Inches values).

    EMU coordinates MUST be integers — PowerPoint rejects decimal values
    in <a:off>/<a:ext> and refuses to open the file. Cast to int defensively.

    ``dash`` accepts any OOXML preset-dash name (e.g., ``"dash"``,
    ``"dashDot"``, ``"sysDash"``).  Default ``None`` = solid line.

    ``head_size`` is the OOXML preset arrowhead size — one of ``'sm'``,
    ``'med'`` (default), or ``'lg'``.  Width and height are set
    together, so passing ``'lg'`` gives a noticeably larger tip while
    leaving the line weight unchanged.
    """
    sx, sy = int(start_xy[0]), int(start_xy[1])
    ex, ey = int(end_xy[0]), int(end_xy[1])
    line = slide.shapes.add_connector(1, sx, sy, ex, ey)  # 1 = STRAIGHT
    line.line.color.rgb = color
    line.line.width = Pt(weight_pt)
    ln = line.line._get_or_add_ln()
    if dash is not None:
        for old in ln.findall(qn('a:prstDash')):
            ln.remove(old)
        prst = ET.SubElement(ln, qn('a:prstDash'))
        prst.set('val', dash)
    if head:
        tailEnd = ET.SubElement(ln, qn('a:tailEnd'))
        tailEnd.set('type', 'triangle')
        tailEnd.set('w', head_size)
        tailEnd.set('h', head_size)
    return line


def _add_wavy_line(slide, x_start, x_end, y_center, *,
                    amplitude=None, cycles=1.75, segments=36,
                    color=NAVY, weight_pt=1.5):
    """Horizontal sinusoidal line from x_start to x_end at y_center.

    Renders a polyline approximation of ``sin(2π · t · cycles)`` inside a
    custGeom shape so the line reads as a gentle wave rather than a
    straight connector.  ``amplitude`` is the peak-to-baseline height in
    EMU; defaults to ~0.04".  ``segments`` controls how finely the wave
    is discretised — 30+ is smooth enough that the polyline reads as a
    curve.  No arrowhead.
    """
    if amplitude is None:
        amplitude = Inches(0.04)
    L = int(x_end - x_start)
    A = int(amplitude)
    if L == 0:
        return None
    flip = "1" if L < 0 else "0"
    bbox_left = int(min(x_start, x_end))
    bbox_top = int(y_center - A)
    bbox_w = abs(L)
    bbox_h = 2 * A

    pts = []
    for i in range(segments + 1):
        t = i / segments
        lx = int(round(t * 100000))
        sin_val = math.sin(2 * math.pi * t * cycles)
        # Path coords: y=0 is top.  Centre at 50000; +1·amplitude → top
        # (0), −1·amplitude → bottom (100000).
        ly = int(round(50000 - sin_val * 50000))
        pts.append((lx, ly))

    path_segs = [f'<a:moveTo><a:pt x="{pts[0][0]}" y="{pts[0][1]}"/></a:moveTo>']
    for lx, ly in pts[1:]:
        path_segs.append(f'<a:lnTo><a:pt x="{lx}" y="{ly}"/></a:lnTo>')
    path_inner = ''.join(path_segs)

    color_hex = f'{color[0]:02X}{color[1]:02X}{color[2]:02X}'
    width_emu = int(weight_pt * 12700)

    P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    A_NS_LOCAL = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    sp_xml = (
        f'<p:sp xmlns:p="{P_NS}" xmlns:a="{A_NS_LOCAL}">'
        f'<p:nvSpPr><p:cNvPr id="0" name="WavyLine"/>'
        f'<p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr>'
        f'<a:xfrm flipH="{flip}">'
        f'<a:off x="{bbox_left}" y="{bbox_top}"/>'
        f'<a:ext cx="{bbox_w}" cy="{bbox_h}"/>'
        f'</a:xfrm>'
        f'<a:custGeom>'
        f'<a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/>'
        f'<a:rect l="0" t="0" r="0" b="0"/>'
        f'<a:pathLst>'
        f'<a:path w="100000" h="100000" fill="none">'
        f'{path_inner}'
        f'</a:path>'
        f'</a:pathLst>'
        f'</a:custGeom>'
        f'<a:noFill/>'
        f'<a:ln w="{width_emu}" cap="rnd">'
        f'<a:solidFill><a:srgbClr val="{color_hex}"/></a:solidFill>'
        f'</a:ln>'
        f'</p:spPr>'
        f'</p:sp>'
    )
    elem = ET.fromstring(sp_xml)
    slide.shapes._spTree.append(elem)
    return elem


def _add_arrow_shape(slide, left, top, width, height, *,
                     direction="right", fill=GOLD, line=None):
    """Block arrow shape (the 'we-are-here' indicator).

    direction: "right" (default), "left", "up", or "down".
    """
    left, top, width, height = int(left), int(top), int(width), int(height)
    geom_map = {
        "left": MSO_SHAPE.LEFT_ARROW,
        "right": MSO_SHAPE.RIGHT_ARROW,
        "up": MSO_SHAPE.UP_ARROW,
        "down": MSO_SHAPE.DOWN_ARROW,
    }
    geom = geom_map.get(direction, MSO_SHAPE.RIGHT_ARROW)
    shp = slide.shapes.add_shape(geom, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp



# --------------------------------------------------------------------------
# Layout 2 — Section Header / Agenda (parameterized)
# --------------------------------------------------------------------------

def make_section_agenda(prs, page_num, *, current_part_idx=None,
                        current_sub_idx=None,
                        section_tag="Module 3 · Agenda",
                        title="Agenda"):
    """Render an agenda / section-divider slide.

    ``current_part_idx`` makes that Part navy and fades the others.
    ``current_sub_idx`` (optional, only meaningful with current_part_idx)
    further restricts the navy highlight inside the current Part to a
    single sub-bullet — the other subs render in faded gray, alongside
    Part 2.  Used by intra-Part dividers (e.g., the new Short Run agenda
    on page 13, or the Long Run agenda on page 31).
    """
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, section_tag)
    _draw_action_title(slide, title)

    y = Inches(1.85)
    part_title_h = Inches(0.6)
    sub_list_h = Inches(1.35)
    block_gap = Inches(0.15)

    for idx, part in enumerate(MODULE_AGENDA):
        if current_part_idx is None:
            color = NAVY  # preview: highlight all Parts
        else:
            color = NAVY if idx == current_part_idx else FADED

        # 2026-05-19: per-sub fade.  When this is the current Part AND a
        # specific sub-index is highlighted, give every other sub the
        # FADED color so the deck shows "you are HERE" within the Part.
        if (current_part_idx is not None and idx == current_part_idx
                and current_sub_idx is not None):
            sub_colors = [
                NAVY if i == current_sub_idx else FADED
                for i in range(len(part["subs"]))
            ]
        else:
            sub_colors = None

        _add_text(slide, MARGIN, y, RULE_W, part_title_h,
                  part["title"], size=30, bold=True, color=color,
                  font="Calibri")
        y += part_title_h

        _add_bulleted_list(
            slide,
            left=MARGIN + Inches(0.4),
            top=y,
            width=RULE_W - Inches(0.4),
            height=sub_list_h,
            items=part["subs"],
            size=24, color=color, bullet_color=color,
            # 2026-05-19: sub-bullets tightened from 10 → 3 pt to match
            # the deck-wide rhythm (sub-bullets cluster under their
            # parent rather than breathing out).  Applied to both
            # Parts on every agenda divider in the deck.
            line_spacing_pts=3,
            autonum_scheme='alphaLcPeriod',
            colors=sub_colors,
        )
        y += sub_list_h + block_gap

    _draw_footer(slide, FOOTER_TEXT, page_num)
    return slide


# --------------------------------------------------------------------------
# Layout 3 — Content bulleted
# --------------------------------------------------------------------------

# --------------------------------------------------------------------------
# Slide titles are set in title case (2026-08-24, Nico), the way a paper
# title is: every significant word starts with a capital, while articles,
# coordinating conjunctions and short prepositions stay lower case unless
# they open the title, close it, or follow a colon.
#
# The pass only ever RAISES a letter - it never lower-cases a word that is
# already capitalised - so acronyms (MR, TR, OLS, WTP, A/B), product names
# and Nico's own capitalisation choices survive untouched.
# --------------------------------------------------------------------------

TITLE_LOWER = {
    "a", "an", "the",
    "and", "but", "or", "nor", "for", "so", "yet",
    "as", "at", "by", "in", "into", "of", "off", "on", "onto", "out",
    "per", "to", "up", "via", "with", "from", "over", "than", "vs",
}


def _tc_word(word, force):
    """Capitalise one whitespace-delimited token, hyphen parts included."""
    parts = word.split("-")
    out = []
    for i, part in enumerate(parts):
        core = part.lstrip("\u201c\u2018\"'([")
        lead = part[:len(part) - len(core)]
        stripped = core.rstrip("\u201d\u2019\"')]:,.?!")
        trail = core[len(stripped):]
        low = stripped.lower().rstrip(".")
        keep_low = (not force) and low in TITLE_LOWER
        if stripped and not keep_low and stripped[0].islower():
            stripped = stripped[0].upper() + stripped[1:]
        out.append(lead + stripped + trail)
        force = False          # only the first hyphen part can be forced
    return "-".join(out)


def _title_case(title):
    words = title.split(" ")
    idx = [i for i, w in enumerate(words) if w.strip()]
    if not idx:
        return title
    first, last = idx[0], idx[-1]
    out = []
    force_next = True
    for i, w in enumerate(words):
        if not w.strip():
            out.append(w)
            continue
        force = force_next or i == first or i == last
        out.append(_tc_word(w, force))
        force_next = w.rstrip().endswith((":", "?", "!", "\u2014", "\u2013"))
    return " ".join(out)


_draw_action_title_raw = _draw_action_title


def _draw_action_title(slide, title, gold_len=GOLD_W):
    return _draw_action_title_raw(slide, _title_case(title), gold_len)


def make_content_bulleted(prs, page_num, section_tag, title, bullets, *,
                          size=24, sub_size=None, line_spacing_pts=18,
                          sub_line_spacing_pts=None,
                          extras=None, bullets_top=None,
                          bullets_width=None):
    """bullets: list of (text, level) tuples OR plain strings (level=0).

    ``bullets_top`` overrides the default body-region start (Inches(1.85))
    — useful when the slide also hosts large diagrams below the bullets
    and the bullets need to lift up to avoid overlap.

    ``sub_line_spacing_pts`` overrides the legacy
    ``max(6, line_spacing_pts - 8)`` formula for sub-bullet space-before.
    """
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, section_tag)
    _draw_action_title(slide, title)

    normalized = [(b, 0) if isinstance(b, str) else b for b in bullets]

    if bullets_top is None:
        bullets_top = Inches(1.85)

    _add_hierarchical_bullets(
        slide,
        left=MARGIN,
        top=bullets_top,
        width=RULE_W if bullets_width is None else bullets_width,
        height=Inches(5.0),
        items=normalized,
        size=size,
        sub_size=sub_size,
        line_spacing_pts=line_spacing_pts,
        sub_line_spacing_pts=sub_line_spacing_pts,
    )

    if extras is not None:
        extras(slide)

    _draw_footer(slide, FOOTER_TEXT, page_num)
    return slide


def _inject_bullet_lst_style(tf, *, size, sub_size,
                              main_color=NAVY, sub_color=GRAY,
                              main_char='▪', sub_char='–',
                              main_space_pts=None, sub_space_pts=None):
    """Inject <a:lstStyle> into a text frame defining per-level defaults.

    Defines lvl1pPr (main bullets, lvl="0") and lvl2pPr (sub bullets,
    lvl="1") with bullet character, indent, color, font, and space-
    before defaults.  Enables PowerPoint's native Tab / Shift+Tab to
    auto-reformat a bullet when its level changes: runs / paragraphs
    that do NOT explicitly override these properties inherit them.

    main_space_pts / sub_space_pts: if set, the <a:spcBef> for each
    level.  Sub-bullets default to a tighter spacing than main bullets
    (3 pt vs ~10 pt) per the deck's pedagogical preference: sub-bullets
    visually cluster under their parent, main bullets give the eye a
    larger break.
    """
    txBody = tf._txBody
    # Remove any prior lstStyle
    for old in txBody.findall(qn('a:lstStyle')):
        txBody.remove(old)
    lstStyle = ET.Element(qn('a:lstStyle'))
    levels = [
        ('a:lvl1pPr', 342900, -342900, main_color, main_char, size, main_space_pts),
        ('a:lvl2pPr', 685800, -228600, sub_color, sub_char, sub_size, sub_space_pts),
    ]
    for tag, mar_l, indent, color, char, sz, space_pts in levels:
        lvl = ET.SubElement(lstStyle, qn(tag))
        lvl.set('marL', str(mar_l))
        lvl.set('indent', str(indent))
        # spcBef must precede bullet attributes (OOXML schema order).
        if space_pts is not None:
            spcBef = ET.SubElement(lvl, qn('a:spcBef'))
            spcPts = ET.SubElement(spcBef, qn('a:spcPts'))
            spcPts.set('val', str(int(space_pts * 100)))
        bc = ET.SubElement(lvl, qn('a:buClr'))
        sc = ET.SubElement(bc, qn('a:srgbClr'))
        sc.set('val', '{:02X}{:02X}{:02X}'.format(color[0], color[1], color[2]))
        bf = ET.SubElement(lvl, qn('a:buFont'))
        bf.set('typeface', 'Calibri')
        bch = ET.SubElement(lvl, qn('a:buChar'))
        bch.set('char', char)
        drp = ET.SubElement(lvl, qn('a:defRPr'))
        drp.set('sz', str(int(sz * 100)))
        drp.set('b', '0')
        sf = ET.SubElement(drp, qn('a:solidFill'))
        sfc = ET.SubElement(sf, qn('a:srgbClr'))
        sfc.set('val', '{:02X}{:02X}{:02X}'.format(color[0], color[1], color[2]))
        lt = ET.SubElement(drp, qn('a:latin'))
        lt.set('typeface', 'Calibri')
    # Insert lstStyle after bodyPr (proper element order in <a:txBody>)
    bodyPr = txBody.find(qn('a:bodyPr'))
    if bodyPr is not None:
        bodyPr.addnext(lstStyle)
    else:
        txBody.insert(0, lstStyle)


def _add_hierarchical_bullets(slide, left, top, width, height, items,
                              *, size=24, sub_size=None, line_spacing_pts=18,
                              sub_line_spacing_pts=None):
    """Render bullets with indent levels.

    Bullet item forms:
        (text, level)                       — simple, defaults from level
        (text, level, opts)                 — opts dict overrides
        (runs_list, level, opts)            — multi-run paragraph

    text:
        - str  → a single run with text
        - ''   → empty paragraph (visual spacer; no run)
        - list → multi-run: list of ``(run_text, run_opts)`` tuples

    Paragraph-level opts (all optional):
        bullet_style: 'main' (▪ NAVY) | 'sub' (– GRAY) | 'arrow' (no bullet,
            plain left-indent; the run text supplies the leader char like
            "→" or Wingdings) | 'none'
        mar_l, indent: bullet positioning (EMU)
        space_before_pts: spcBef in pts (overrides legacy formula)
        size, color, bold, italic: defaults applied to every run unless
            the run_opts override them.
        align: a PP_ALIGN value (default LEFT).

    Run-level opts (run_opts in a runs_list tuple):
        font_name (default 'Calibri'), size, color, bold, italic,
        underline (bool), wingdings (bool — emits <a:sym typeface="Wingdings"/>
        so a private-use-area character renders as its Wingdings glyph).

    Each paragraph also receives a ``lvl="N"`` attribute when level > 0
    so PowerPoint's Tab / Shift-Tab outline navigation can find an
    explicit outline level.
    """
    if sub_size is None:
        sub_size = size - 4

    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    # 2026-05-18: inject <a:lstStyle> with per-level defaults so that
    # PowerPoint's Tab / Shift+Tab keyboard shortcuts auto-reformat the
    # bullet when its outline level changes.  Default 'main'/'sub'
    # bullets (no custom indent, level < 2) skip explicit pPr/run
    # styling — they inherit from lstStyle.  Bullets with custom indent
    # OR level >= 2 OR explicit size/color overrides still apply
    # explicit styling.
    # 2026-05-18 (later): also moved space-before into lstStyle so a
    # paragraph demoted via Tab picks up the sub-bullet's tighter spcBef
    # automatically.  Sub-bullet default: 3 pt (was 6 pt) — per user
    # preference, sub-bullets cluster tightly under their parent.
    sub_default_space = (sub_line_spacing_pts if sub_line_spacing_pts is not None
                          else 3)
    _inject_bullet_lst_style(tf, size=size, sub_size=sub_size,
                              main_space_pts=line_spacing_pts,
                              sub_space_pts=sub_default_space)

    for i, item in enumerate(items):
        if len(item) == 2:
            text, level = item
            opts = {}
        else:
            text, level, opts = item
            opts = opts or {}

        # Determine inheritance up-front so spcBef logic below can use it.
        style = opts.get('bullet_style', 'main' if level == 0 else 'sub')
        has_custom_indent = 'mar_l' in opts or 'indent' in opts
        inherits_from_lst = (style in ('main', 'sub')
                              and not has_custom_indent
                              and level < 2)

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        # 2026-08-25: paragraphs can be centred (the video deck's
        # "compute:" list on slide 42 is a centred, un-bulleted list)
        p.alignment = opts.get('align', PP_ALIGN.LEFT)
        pPr = p._p.get_or_add_pPr()

        # Space-before.  When the paragraph inherits from lstStyle and the
        # caller has not overridden, skip the per-pPr spcBef so that Tab-
        # induced level changes pick up the new level's spcBef from
        # lstStyle.  Otherwise set explicitly.
        if i > 0:
            sp_override = opts.get('space_before_pts')
            if sp_override is not None:
                spcBef = ET.SubElement(pPr, qn('a:spcBef'))
                pts = ET.SubElement(spcBef, qn('a:spcPts'))
                pts.set('val', str(sp_override * 100))
            elif not inherits_from_lst:
                if level == 0:
                    sp = line_spacing_pts
                else:
                    sp = sub_default_space
                spcBef = ET.SubElement(pPr, qn('a:spcBef'))
                pts = ET.SubElement(spcBef, qn('a:spcPts'))
                pts.set('val', str(sp * 100))
            # else: inherit from lstStyle.

        # Outline-level attribute (enables PowerPoint Tab/Shift-Tab)
        if level > 0:
            pPr.set('lvl', str(level))

        # Bullet styling — default 'main' (lvl 0) and 'sub' (lvl 1)
        # inherit from lstStyle when no custom indent is requested.
        # 'arrow' and 'none' suppress the lstStyle bullet via <a:buNone/>.
        # Level >= 2 still uses explicit _set_bullet_char (lstStyle here
        # only defines lvl1pPr and lvl2pPr).
        if style == 'main':
            if not inherits_from_lst:
                _set_bullet_char(p, char='▪', color=NAVY,
                                  mar_l=opts.get('mar_l', 342900),
                                  indent=opts.get('indent', -342900),
                                  size_pct=100)
        elif style == 'sub':
            if not inherits_from_lst:
                default_mar = 342900 + level * 342900
                _set_bullet_char(p, char='–', color=GRAY,
                                  mar_l=opts.get('mar_l', default_mar),
                                  indent=opts.get('indent', -228600),
                                  size_pct=100)
        elif style == 'arrow':
            # Plain left-indent, no bullet glyph; user text supplies leader.
            pPr.set('marL', str(opts.get('mar_l', 457200)))
            if 'indent' in opts:
                pPr.set('indent', str(opts['indent']))
            # Explicit <a:buNone/> so the lstStyle bullet doesn't bleed
            # through.
            ET.SubElement(pPr, qn('a:buNone'))
        elif style == 'none':
            ET.SubElement(pPr, qn('a:buNone'))

        # Empty paragraph (spacer) — no run; suppress the bullet so the
        # empty line doesn't render an orphan glyph.
        if text == '':
            if inherits_from_lst:
                ET.SubElement(pPr, qn('a:buNone'))
            continue

        # Normalize text to a runs list
        if isinstance(text, str):
            runs = [(text, {})]
        else:
            runs = text  # already a list of (run_text, run_opts) tuples

        # Paragraph-level explicit overrides (None = inherit from lstStyle).
        para_size_override = opts.get('size')
        para_color_override = opts.get('color')
        para_bold_override = opts.get('bold')
        para_italic_override = opts.get('italic')

        for run_text, run_opts in runs:
            run_opts = run_opts or {}

            # Inline OMML math zone — append <a14:m><m:oMath>...</m:oMath></a14:m>
            # as a sibling of the regular <a:r> runs so the formula renders
            # in-line with surrounding prose.  ``run_text`` is interpreted
            # as raw OMML content (e.g., from `_formula_mp_ratio` or the
            # `_omml_*` builders).
            if run_opts.get('omml'):
                om_sz = run_opts.get('size')
                if om_sz is None:
                    om_sz = para_size_override
                if om_sz is None:
                    om_sz = size if level == 0 else sub_size
                om_clr = run_opts.get('color') or para_color_override
                if om_clr is None:
                    om_clr = NAVY if level == 0 else GRAY
                sz_centi = int(om_sz * 100)
                clr_hex = '{:02X}{:02X}{:02X}'.format(
                    om_clr[0], om_clr[1], om_clr[2])
                a14_xml = (
                    f'<a14:m xmlns:a14="{A14_NS}" '
                    f'xmlns:m="{M_NS}" xmlns:a="{A_NS}">'
                    f'<m:oMath>{run_text}</m:oMath>'
                    f'</a14:m>'
                )
                a14_elem = ET.fromstring(a14_xml)
                for r_elem in a14_elem.iter(qn('m:r')):
                    arPr = r_elem.find(qn('a:rPr'))
                    if arPr is None:
                        arPr = ET.Element(qn('a:rPr'))
                        r_elem.insert(0, arPr)
                    arPr.set('sz', str(sz_centi))
                    if arPr.get('lang') is None:
                        arPr.set('lang', 'en-US')
                    for sf in arPr.findall(qn('a:solidFill')):
                        arPr.remove(sf)
                    # fill must come BEFORE a:latin (schema order), else
                    # PowerPoint ignores the color
                    sf = arPr.makeelement(qn('a:solidFill'), {})
                    srgb = ET.SubElement(sf, qn('a:srgbClr'))
                    srgb.set('val', clr_hex)
                    arPr.insert(0, sf)
                p._p.append(a14_elem)
                continue

            run = p.add_run()
            run.text = run_text
            run.font.name = run_opts.get('font_name', 'Calibri')

            # Size: run-opt > para-opt > inherit (for inherits_from_lst
            # cases) or fall back to level default (for explicit-styled
            # paragraphs).
            run_sz = run_opts.get('size')
            sz = run_sz if run_sz is not None else para_size_override
            if sz is None and not inherits_from_lst:
                sz = size if level == 0 else sub_size
            if sz is not None:
                run.font.size = Pt(sz)

            run_clr = run_opts.get('color')
            clr = run_clr if run_clr is not None else para_color_override
            if clr is None and not inherits_from_lst:
                clr = NAVY if level == 0 else GRAY
            if clr is not None:
                run.font.color.rgb = clr

            # Bold: only set if explicitly requested.  Default = inherit
            # (lstStyle defRPr has b="0").
            run_b = run_opts.get('bold')
            b = run_b if run_b is not None else para_bold_override
            if b is not None:
                run.font.bold = b

            it = run_opts.get('italic', para_italic_override)
            if it is not None:
                run.font.italic = it
            if run_opts.get('underline'):
                run.font.underline = True
            if run_opts.get('highlight'):
                # 2026-08-25: <a:highlight> is a background wash behind
                # the glyphs, not a font colour - Nico marks the "%" on
                # slide 37 this way
                rPr = run._r.find(qn('a:rPr'))
                if rPr is None:
                    rPr = run._r.makeelement(qn('a:rPr'), {})
                    run._r.insert(0, rPr)
                hl = ET.SubElement(rPr, qn('a:highlight'))
                clr = ET.SubElement(hl, qn('a:srgbClr'))
                clr.set('val', run_opts['highlight'])
            if run_opts.get('wingdings'):
                # Add <a:sym typeface="Wingdings"/> so private-use-area
                # characters render as their Wingdings glyphs.
                rPr = run._r.find(qn('a:rPr'))
                if rPr is None:
                    rPr = run._r.makeelement(qn('a:rPr'), {})
                    run._r.insert(0, rPr)
                sym = ET.SubElement(rPr, qn('a:sym'))
                sym.set('typeface', 'Wingdings')

    return box


# --------------------------------------------------------------------------
# Layout-3 variant — diagram canvas (action title + free-form shapes below).
# Used for the agenda flowchart and the Big Picture diagram so they live
# inside the same visual chrome as content slides but render a diagram
# instead of bullets.
# --------------------------------------------------------------------------

def make_diagram_slide(prs, page_num, section_tag, title, draw_diagram):
    """Action-title slide with free-form diagram region below.

    draw_diagram: callable(slide) that renders the diagram in the body
    region (approximately y = 1.85 to 6.95).
    """
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, section_tag)
    _draw_action_title(slide, title)
    draw_diagram(slide)
    _draw_footer(slide, FOOTER_TEXT, page_num)
    return slide


# --------------------------------------------------------------------------
# Layout 5 — Poll slide (A./B./C./D. options + POLL pill, no QR box)
# --------------------------------------------------------------------------

def _draw_poll_pill(slide, *, position='top-right',
                     fill=NAVY, text_color=WHITE, dot_color=GOLD,
                     width=None, height=None, text_size=14,
                     shadow=False):
    """Small 'POLL' pill, right-aligned.

    position: 'top-right' (default, just below the top bar) or
              'bottom-right' (bottom band aligned with discussion break).
    width / height: override the default 1.05" × 0.42" pill size.
    text_size: "POLL" font size (pt).  Scales up when the pill is enlarged.
    dot_color: pass None to skip the leading accent dot.
    shadow: add a soft drop shadow to the pill (matches discussion-break
        chrome).  Default False to preserve the original flat top-right
        pill look.
    """
    pill_w = width if width is not None else Inches(1.05)
    pill_h = height if height is not None else Inches(0.42)
    pill_x = SLIDE_W - MARGIN - pill_w
    if position == 'bottom-right':
        # Bottom-aligned with where _add_discussion_break ends
        # (top=6.25, height=0.72 → bottom=6.97).  Pill height varies, so
        # pick top so the pill BOTTOM lands on the same 6.97 baseline.
        pill_y = Inches(6.97) - pill_h
    else:
        pill_y = Inches(0.55)

    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  pill_x, pill_y, pill_w, pill_h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    if shadow:
        _add_drop_shadow(shp)
    else:
        shp.shadow.inherit = False
    try:
        shp.adjustments[0] = 0.5
    except Exception:
        pass
    _add_text(slide, pill_x, pill_y, pill_w, pill_h,
              "POLL", size=text_size, bold=True, color=text_color, font="Calibri",
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if dot_color is not None:
        dot_size = Inches(0.16)
        dot_x = pill_x - Inches(0.16) - dot_size
        dot_y = pill_y + (pill_h - dot_size) // 2
        _add_rect(slide, dot_x, dot_y, dot_size, dot_size, dot_color)


def make_poll_slide(prs, page_num, section_tag, title, options, *,
                    instructions="Respond at PollEv.com/nvoigtlaender",
                    size=30, line_spacing_pts=22):
    """Layout 5 – Poll slide.

    options: list of strings (A./B./C./D. auto-numbering applied).
    """
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, section_tag)
    _draw_action_title(slide, title)
    _draw_poll_pill(slide)

    _add_bulleted_list(
        slide,
        left=MARGIN,
        top=Inches(2.0),
        width=RULE_W,
        height=Inches(4.4),
        items=options,
        size=size, color=NAVY, bullet_color=NAVY,
        line_spacing_pts=line_spacing_pts,
        autonum_scheme='alphaUcPeriod',
    )

    _add_text(slide, MARGIN, Inches(6.55), RULE_W, Inches(0.4),
              instructions, size=16, italic=True, color=GRAY,
              align=PP_ALIGN.RIGHT)

    _draw_footer(slide, FOOTER_TEXT, page_num)
    return slide


# --------------------------------------------------------------------------
# Slide-jump hyperlinks — make a run or shape clickable in slideshow
# mode so it advances PowerPoint to a different slide.  python-pptx
# doesn't expose this directly, so we manipulate the rPr / cNvPr XML
# and register the slide relationship via the public OPC API.
# --------------------------------------------------------------------------

def _add_slide_jump_hyperlink_run(source_slide, run, target_slide,
                                    *, lock_color=True, underline=True):
    """Make ``run`` a clickable hyperlink that advances to ``target_slide``.

    ``lock_color=True`` adds the Office hyperlinkcolor extension so the
    run's explicit ``solidFill`` color is preserved at render time —
    without it, PowerPoint repaints hyperlink text in the theme's
    ``hlink`` color (a light blue) regardless of what the rPr says.

    ``underline=True`` sets ``u="sng"`` so the hyperlink stays visibly
    underlined even when the color-lock suppresses the theme styling.
    """
    rId = source_slide.part.relate_to(target_slide.part, RT.SLIDE)
    rPr = run._r.get_or_add_rPr()
    if underline:
        rPr.set('u', 'sng')
    for hl in rPr.findall(qn('a:hlinkClick')):
        rPr.remove(hl)
    hlinkClick = ET.SubElement(rPr, qn('a:hlinkClick'))
    hlinkClick.set(qn('r:id'), rId)
    hlinkClick.set('action', 'ppaction://hlinksldjump')
    if lock_color:
        AHYP_NS = 'http://schemas.microsoft.com/office/drawing/2018/hyperlinkcolor'
        ext_xml = (
            f'<a:extLst xmlns:a="{A_NS}">'
            f'<a:ext uri="{{A12FA001-AC4F-418D-AE19-62706E023703}}">'
            f'<ahyp:hlinkClr xmlns:ahyp="{AHYP_NS}" val="tx"/>'
            f'</a:ext>'
            f'</a:extLst>'
        )
        hlinkClick.append(ET.fromstring(ext_xml))


def _add_slide_jump_hyperlink_shape(source_slide, shape, target_slide):
    """Make ``shape`` clickable so the whole shape jumps to ``target_slide``."""
    rId = source_slide.part.relate_to(target_slide.part, RT.SLIDE)
    nvSpPr = shape._element.find(qn('p:nvSpPr'))
    cNvPr = nvSpPr.find(qn('p:cNvPr')) if nvSpPr is not None else None
    if cNvPr is None:
        return
    for hl in cNvPr.findall(qn('a:hlinkClick')):
        cNvPr.remove(hl)
    hlinkClick = ET.SubElement(cNvPr, qn('a:hlinkClick'))
    hlinkClick.set(qn('r:id'), rId)
    hlinkClick.set('action', 'ppaction://hlinksldjump')


def _link_shape_to_url(slide, shape, url, *, tooltip=None):
    """Make a shape or picture open an external URL when clicked.

    2026-08-26: restores the podcast link Nico had on the original
    Uber slide.  A CLICK link rather than a hover link — a hover link
    would jump to the browser whenever the mouse crossed the box during
    a lecture; the ScreenTip gives the hover feedback instead.
    """
    shape.click_action.hyperlink.address = url
    cNvPr = shape._element.find('.//' + qn('p:cNvPr'))
    if tooltip and cNvPr is not None:
        hl = cNvPr.find(qn('a:hlinkClick'))
        if hl is not None:
            hl.set('tooltip', tooltip)
    return shape


# --------------------------------------------------------------------------
# Image helpers — embed source-deck images into the new deck.
# Images are pre-extracted to _source_images/slide{N}_{rId}.{ext}.
# --------------------------------------------------------------------------

SRC_IMG_DIR = Path(__file__).parent / "_source_images"


def _add_drop_shadow(shape, *, blur="50800", dist="38100",
                      direction="2700000", alpha="45000"):
    """Add a soft drop shadow to any shape that exposes spPr (pictures,
    rectangles, rounded rects).  Default: 4pt blur, 3pt offset, 45° down-
    right, 45% opacity black.  Used deck-wide for figures/boxes."""
    try:
        spPr = shape._element.spPr
    except AttributeError:
        # Fallback for shapes whose XML element exposes spPr only via find()
        spPr = shape._element.find(qn('p:spPr'))
    if spPr is None:
        return shape
    for old in spPr.findall(qn('a:effectLst')):
        spPr.remove(old)
    effLst = ET.SubElement(spPr, qn('a:effectLst'))
    outerShdw = ET.SubElement(effLst, qn('a:outerShdw'))
    outerShdw.set('blurRad', str(blur))
    outerShdw.set('dist', str(dist))
    outerShdw.set('dir', str(direction))
    outerShdw.set('algn', 'tl')
    outerShdw.set('rotWithShape', '0')
    rgb = ET.SubElement(outerShdw, qn('a:srgbClr'))
    rgb.set('val', '000000')
    a = ET.SubElement(rgb, qn('a:alpha'))
    a.set('val', str(alpha))
    return shape


def _add_graphicframe_shadow(slide, left, top, width, height, *,
                              shadow_alpha=45000, rounded=False,
                              corner_in=0.16):
    """White backing rectangle with an outerShdw effect, behind a
    graphicFrame (table or chart).  graphicFrames can't host
    a:effectLst directly; this rect supplies the shadow projected
    OUTSIDE its bounds.

    Charts have transparent plot areas by default, so a coloured backing
    bleeds through and tints the figure.  Use white instead — the chart
    sees white through its own transparent areas (clean background) and
    the shadow renders only at the visible edges.  Call BEFORE adding
    the table/chart so z-order is correct.
    """
    shdw = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        int(left), int(top), int(width), int(height),
    )
    if rounded:
        # adj is a fraction of half the SHORTER side - fix the radius in
        # inches so a wide card is not over-rounded
        try:
            shdw.adjustments[0] = min(
                0.5, corner_in / (min(width, height) / 914400.0 / 2.0))
        except Exception:
            pass
    shdw.fill.solid()
    shdw.fill.fore_color.rgb = WHITE
    shdw.line.fill.background()
    shdw.shadow.inherit = False
    sp_pr = shdw._element.spPr
    # Strip any default <a:effectLst/> python-pptx may have inserted
    # (duplicate effectLst makes PowerPoint refuse to open the file).
    for old in sp_pr.findall(qn('a:effectLst')):
        sp_pr.remove(old)
    effLst = ET.SubElement(sp_pr, qn('a:effectLst'))
    outerShdw = ET.SubElement(effLst, qn('a:outerShdw'))
    outerShdw.set('blurRad', '50800')
    outerShdw.set('dist', '38100')
    outerShdw.set('dir', '2700000')
    outerShdw.set('algn', 'tl')
    outerShdw.set('rotWithShape', '0')
    rgb = ET.SubElement(outerShdw, qn('a:srgbClr'))
    rgb.set('val', '000000')
    a = ET.SubElement(rgb, qn('a:alpha'))
    a.set('val', str(int(shadow_alpha)))
    return shdw


def _add_source_image(slide, src_slide_no, rid, *, left, top, width=None,
                      height=None, shadow=True, rounded=False):
    """Place a source-deck image on the new slide.

    `shadow=True` (default) adds a soft drop shadow so figures pop off the
    background — applied deck-wide per the latest visual direction.  Set
    `shadow=False` for niche cases (transparent PNGs, screenshots that
    already include a shadow, etc.).

    `rounded=True` gives real photographs the deck-standard rounded corners
    (routes through :func:`_apply_picture_style`, which sets both the
    rounded geometry and the drop shadow).  Leave False for logos,
    screenshots, and framed images per the Pictures guideline.
    """
    candidates = list(SRC_IMG_DIR.glob(f"slide{src_slide_no}_{rid}.*"))
    if not candidates:
        return None
    img = candidates[0]
    kwargs = {"left": int(left), "top": int(top)}
    if width is not None:
        kwargs["width"] = int(width)
    if height is not None:
        kwargs["height"] = int(height)
    pic = slide.shapes.add_picture(str(img), **kwargs)
    if rounded:
        _apply_picture_style(pic, corner_pct=6)
    elif shadow:
        _add_drop_shadow(pic)
    return pic


def _apply_picture_style(pic, *, corner_pct=8,
                          shadow_blur=50800, shadow_dist=38100,
                          shadow_dir=2700000, shadow_alpha=50000):
    """Apply rounded corners + drop shadow to a picture shape.

    corner_pct: rounded-corner radius as percent of shorter side (8 ≈ subtle).
    shadow_blur/dist: EMU; defaults give a soft 4pt blur, 3pt offset.
    shadow_dir: 2700000 = 45° down-right (standard).
    shadow_alpha: 50000 = 50% opacity black shadow.
    """
    spPr = pic._element.find(qn('p:spPr'))
    if spPr is None:
        return pic
    # Replace any existing prstGeom with roundRect at corner_pct
    for old in spPr.findall(qn('a:prstGeom')):
        spPr.remove(old)
    prstGeom = ET.Element(qn('a:prstGeom'))
    prstGeom.set('prst', 'roundRect')
    avLst = ET.SubElement(prstGeom, qn('a:avLst'))
    gd = ET.SubElement(avLst, qn('a:gd'))
    gd.set('name', 'adj')
    gd.set('fmla', f'val {int(corner_pct * 1000)}')
    # prstGeom must come after a:xfrm
    xfrm = spPr.find(qn('a:xfrm'))
    if xfrm is not None:
        xfrm.addnext(prstGeom)
    else:
        spPr.insert(0, prstGeom)
    # Replace any existing effectLst with a single outer shadow
    for old in spPr.findall(qn('a:effectLst')):
        spPr.remove(old)
    effectLst = ET.SubElement(spPr, qn('a:effectLst'))
    outerShdw = ET.SubElement(effectLst, qn('a:outerShdw'))
    outerShdw.set('blurRad', str(int(shadow_blur)))
    outerShdw.set('dist', str(int(shadow_dist)))
    outerShdw.set('dir', str(int(shadow_dir)))
    outerShdw.set('algn', 'tl')
    outerShdw.set('rotWithShape', '0')
    rgb = ET.SubElement(outerShdw, qn('a:srgbClr'))
    rgb.set('val', '000000')
    alpha = ET.SubElement(rgb, qn('a:alpha'))
    alpha.set('val', str(int(shadow_alpha)))
    return pic


# --------------------------------------------------------------------------
# Illustrative callout boxes — the "major-concept" badges, "Teaching Note"
# bars, and bottom-of-slide takeaway bands that recur throughout the source
# deck.  Replicating them faithfully (in template colors) is what makes the
# slides feel like the original, not a stripped-down rewrite.
# --------------------------------------------------------------------------

def _set_wingdings(run):
    """Tag a run so its private-use characters render as Wingdings."""
    rPr = run._r.find(qn('a:rPr'))
    if rPr is None:
        rPr = run._r.makeelement(qn('a:rPr'), {})
        run._r.insert(0, rPr)
    sym = ET.SubElement(rPr, qn('a:sym'))
    sym.set('typeface', 'Wingdings')
    return run


def _add_takeaway_bar(slide, text, *, top=Inches(6.4), width=None,
                       height=Inches(0.55), left=None,
                       fill=GOLD, text_color=WHITE,
                       size=20, font="Calibri", bold=True,
                       rounded=False, shadow=False, wingding_lead=None):
    """Bottom-of-slide takeaway band — the 'major concept' callout.

    rounded=True  → ROUNDED_RECTANGLE with ~30% corner radius.
    shadow=True   → soft drop shadow (off by default to keep the flat
                    look on the majority of slides).
    left=None     → horizontally centered. Pass an EMU/Inches value to
                    pin the bar's left edge (used for hand-positioned
                    takeaways).
    """
    if width is None:
        width = Inches(9.6)
    if left is None:
        left = (SLIDE_W - width) // 2
    if not (rounded or shadow):
        return _add_filled_box(slide, left, top, width, height, text,
                                fill=fill, text_color=text_color,
                                size=size, bold=bold, font=font)
    # Inline build for the rounded / shadowed variant
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, int(left), int(top), int(width), int(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    if rounded:
        try: shape.adjustments[0] = 0.30
        except Exception: pass
    shape.shadow.inherit = False
    if shadow:
        _add_drop_shadow(shape)
    tf = shape.text_frame
    tf.margin_left = tf.margin_right = Inches(0.1)
    tf.margin_top = tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    if wingding_lead:
        # 2026-08-25: Nico types the arrow as a Wingdings character, so
        # it needs its own run with <a:sym typeface="Wingdings"/>
        lead = p.add_run()
        lead.text = wingding_lead
        lead.font.size = Pt(size)
        lead.font.bold = bold
        lead.font.color.rgb = text_color
        _set_wingdings(lead)
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = text_color
    return shape


def _add_teaching_note(slide, text, *, top=Inches(6.6), width=None,
                        height=Inches(0.6), left=None,
                        rounded=False, pdf_icon=False, label_color=GOLD,
                        fill_rgb=None):
    """External-document reference card.

    Visually distinct from in-slide callouts: cream/parchment fill, navy
    DASHED border, italic navy text, with a small page-icon glyph on the
    left — clearly signals 'this links to an external Teaching Note doc'.

    Optional styling for the slide-32 (page 33) treatment:
      • ``left``       — pin to an absolute X (default: horizontally center)
      • ``rounded``    — rounded corners + soft drop shadow (default flat)
      • ``pdf_icon``   — render the page glyph as a white folded-corner
                        with bold "PDF" text, sized to nearly fill the card
                        (default: small navy folded-corner, no text)
      • ``label_color``— color of the "SEE TEACHING NOTE →" prefix
                        (default GOLD; pass NAVY for the page-33 look).
      • ``fill_rgb``   — override the cream/parchment card fill
                        (default RGBColor(0xF4, 0xF1, 0xEA)).
    """
    if width is None:
        width = Inches(8.0)
    if left is None:
        left = int((SLIDE_W - width) // 2)
    left, top, width, height = int(left), int(top), int(width), int(height)

    # Card: cream fill, dashed navy border.  ``rounded=True`` switches the
    # base shape to ROUNDED_RECTANGLE and adds a soft drop shadow.
    base_shape = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    card = slide.shapes.add_shape(base_shape, left, top, width, height)
    if rounded:
        try: card.adjustments[0] = 0.20
        except Exception: pass
    card.fill.solid()
    card.fill.fore_color.rgb = (fill_rgb if fill_rgb is not None
                                  else RGBColor(0xF4, 0xF1, 0xEA))  # parchment cream default
    card.line.color.rgb = NAVY
    card.line.width = Pt(1.25)
    card.shadow.inherit = False
    if rounded:
        _add_drop_shadow(card)
    # Dashed line style
    ln = card.line._get_or_add_ln()
    # Remove any existing prstDash
    for el in ln.findall(qn('a:prstDash')):
        ln.remove(el)
    prstDash = ET.SubElement(ln, qn('a:prstDash'))
    prstDash.set('val', 'dash')

    # Empty text on the card; we add the icon + text as separate textboxes
    tf = card.text_frame
    tf.text = ""

    # Small "page" icon on the left – a folded-corner shape.  In default
    # mode it's a small filled-navy glyph (logo-ish); in pdf_icon mode it's
    # a larger WHITE sheet with thin navy border + bold "PDF" text so it
    # reads as a generic PDF document at a glance.
    icon_size = Inches(0.5) if pdf_icon else Inches(0.4)
    icon_x = left + Inches(0.2)
    icon_y = top + (height - icon_size) // 2
    page = slide.shapes.add_shape(MSO_SHAPE.FOLDED_CORNER,
                                   int(icon_x), int(icon_y),
                                   int(icon_size), int(icon_size))
    page.fill.solid()
    if pdf_icon:
        page.fill.fore_color.rgb = WHITE
        page.line.color.rgb = NAVY
        page.line.width = Pt(0.75)
    else:
        page.fill.fore_color.rgb = NAVY
        page.line.fill.background()
    page.shadow.inherit = False
    if pdf_icon:
        # "PDF" label inside the white sheet.
        ptf = page.text_frame
        ptf.word_wrap = False
        ptf.margin_left = 0
        ptf.margin_right = 0
        ptf.margin_top = 0
        ptf.margin_bottom = 0
        ptf.vertical_anchor = MSO_ANCHOR.MIDDLE
        pp = ptf.paragraphs[0]
        pp.alignment = PP_ALIGN.CENTER
        pr = pp.add_run()
        pr.text = "PDF"
        pr.font.name = "Calibri"
        pr.font.size = Pt(11)
        pr.font.bold = True
        pr.font.color.rgb = NAVY

    # Label text (italic, navy) inside the card to the right of the icon
    txt_left = int(icon_x + icon_size + Inches(0.2))
    txt_w = int(width - (icon_x + icon_size + Inches(0.4) - left))
    label_box = slide.shapes.add_textbox(txt_left, top,
                                          txt_w, height)
    label_tf = label_box.text_frame
    label_tf.word_wrap = True
    label_tf.margin_left = 0
    label_tf.margin_right = 0
    label_tf.margin_top = 0
    label_tf.margin_bottom = 0
    label_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = label_tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    # First run: "See teaching note:" prefix in caller-supplied color
    # (GOLD by default; NAVY for the slide-33 treatment).
    r1 = p.add_run()
    r1.text = "SEE TEACHING NOTE  →  "
    r1.font.name = "Calibri"
    r1.font.size = Pt(12)
    r1.font.bold = True
    r1.font.color.rgb = label_color
    # Second run: the title of the note, italic navy
    r2 = p.add_run()
    r2.text = text
    r2.font.name = "Calibri"
    r2.font.size = Pt(16)
    r2.font.italic = True
    r2.font.bold = True
    r2.font.color.rgb = NAVY
    return card


def _add_discussion_break(slide, *, top=Inches(6.25), width=Inches(4.8),
                           left=None, text="Discussion Break"):
    """Rounded-parallelogram 'discussion break' badge (bottom-right).

    Custom-geometry shape: top and bottom edges are horizontal; the left
    and right edges slant at 45° in real space (skew = height of the
    shape).  All four corners are slightly rounded.  Gold fill, navy
    bold text, soft drop shadow.

    left=None → pin to the right edge with the default MARGIN. Pass an
    EMU/Inches value to override (used for hand-positioned badges).
    """
    height = Inches(0.72)
    if left is None:
        left = SLIDE_W - MARGIN - width
    left, top, width, height = int(left), int(top), int(width), int(height)
    # Compute skew in path-coordinate units so that left/right sides slant
    # at 45° in REAL space:  horizontal offset of top edge == shape height.
    skew = int(100000 * height / width) if width else 15000
    skew = min(max(skew, 6000), 35000)
    r = 5000          # corner radius in path units — gentle rounding

    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = GOLD
    shp.line.fill.background()
    shp.shadow.inherit = False
    # Strip the default prstGeom – we'll inject custGeom instead.
    spPr = shp._element.spPr
    for old in spPr.findall(qn('a:prstGeom')):
        spPr.remove(old)
    rs = int(r * skew / 100000)
    # IMPORTANT: <a:rect> defines the TEXT bounding rectangle inside the
    # custom geometry.  Set it to the parallelogram's inscribed rectangle
    # (from TL vertex to BR vertex) so PowerPoint won't render text past
    # the slanted edges, regardless of the text frame's own lIns/rIns.
    custgeom_xml = (
        f'<a:custGeom xmlns:a="{A_NS}">'
        f'<a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/>'
        f'<a:rect l="{skew}" t="0" r="{100000-skew}" b="100000"/>'
        f'<a:pathLst><a:path w="100000" h="100000">'
        # Top-left rounded corner (vertex at (skew, 0))
        f'<a:moveTo><a:pt x="{skew+r}" y="0"/></a:moveTo>'
        # Top edge
        f'<a:lnTo><a:pt x="{100000-r}" y="0"/></a:lnTo>'
        # Top-right corner round
        f'<a:cubicBezTo>'
        f'<a:pt x="100000" y="0"/><a:pt x="100000" y="0"/>'
        f'<a:pt x="{100000-rs}" y="{r}"/>'
        f'</a:cubicBezTo>'
        # Right slanted side (down-left)
        f'<a:lnTo><a:pt x="{100000-skew+rs}" y="{100000-r}"/></a:lnTo>'
        # Bottom-right corner round (vertex at (100000-skew, 100000))
        f'<a:cubicBezTo>'
        f'<a:pt x="{100000-skew}" y="100000"/>'
        f'<a:pt x="{100000-skew}" y="100000"/>'
        f'<a:pt x="{100000-skew-r}" y="100000"/>'
        f'</a:cubicBezTo>'
        # Bottom edge
        f'<a:lnTo><a:pt x="{r}" y="100000"/></a:lnTo>'
        # Bottom-left corner round (vertex at (0, 100000))
        f'<a:cubicBezTo>'
        f'<a:pt x="0" y="100000"/><a:pt x="0" y="100000"/>'
        f'<a:pt x="{rs}" y="{100000-r}"/>'
        f'</a:cubicBezTo>'
        # Left slanted side (up-right)
        f'<a:lnTo><a:pt x="{skew-rs}" y="{r}"/></a:lnTo>'
        # Top-left corner round (close the path)
        f'<a:cubicBezTo>'
        f'<a:pt x="{skew}" y="0"/><a:pt x="{skew}" y="0"/>'
        f'<a:pt x="{skew+r}" y="0"/>'
        f'</a:cubicBezTo>'
        f'<a:close/>'
        f'</a:path></a:pathLst>'
        f'</a:custGeom>'
    )
    custgeom = ET.fromstring(custgeom_xml)
    # Insert custGeom right after a:xfrm (schema order)
    xfrm = spPr.find(qn('a:xfrm'))
    if xfrm is not None:
        xfrm.addnext(custgeom)
    else:
        spPr.insert(0, custgeom)

    # Drop shadow (45° down-right, 50% opacity).
    for old in spPr.findall(qn('a:effectLst')):
        spPr.remove(old)
    effectLst = ET.SubElement(spPr, qn('a:effectLst'))
    outerShdw = ET.SubElement(effectLst, qn('a:outerShdw'))
    outerShdw.set('blurRad', '50800')
    outerShdw.set('dist', '38100')
    outerShdw.set('dir', '2700000')
    outerShdw.set('algn', 'tl')
    outerShdw.set('rotWithShape', '0')
    rgb = ET.SubElement(outerShdw, qn('a:srgbClr'))
    rgb.set('val', '000000')
    alpha = ET.SubElement(rgb, qn('a:alpha'))
    alpha.set('val', '50000')

    # Text — render as a SEPARATE textbox overlaid on top of the
    # parallelogram, positioned exactly inside the inscribed rectangle.
    # This decouples text placement from the shape geometry and avoids
    # PowerPoint rendering the run past the slanted edges (which can
    # happen with the in-shape text frame on some PowerPoint versions
    # even with <a:rect> set inside custGeom).
    # In real EMU, the skew equals the shape height (45° slant), so the
    # inscribed rectangle spans (left + height) → (left + width - height).
    skew_emu = height
    ins_left = left + skew_emu
    ins_top = top
    ins_w = width - 2 * skew_emu
    ins_h = height
    txt = slide.shapes.add_textbox(int(ins_left), int(ins_top),
                                     int(ins_w), int(ins_h))
    tf = txt.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(28)         # bumped 20 → 28 on 2026-05-16
    run.font.bold = True
    run.font.color.rgb = NAVY
    return shp


def _add_callout_box(slide, left, top, width, height, text, *,
                      fill=GOLD, text_color=WHITE, size=14, bold=True,
                      rounded=False, shadow=False):
    """Small free-form annotation/callout (e.g., 'plot the slope', 'Revenue
    per car net of material cost').  Used to mark a graph or sub-region.

    ``rounded``/``shadow`` (default off) route to the rounded + drop-shadow
    variant for the deck-wide box treatment."""
    if rounded or shadow:
        return _add_rounded_filled_box(
            slide, left, top, width, height, text,
            fill=fill, text_color=text_color,
            size=size, bold=bold, font="Calibri",
            corner_pct=0.12, shadow=shadow)
    return _add_filled_box(slide, left, top, width, height, text,
                            fill=fill, text_color=text_color,
                            size=size, bold=bold, font="Calibri")


def _add_anchor_burst(slide, left, top, width, height,
                       top_text, bottom_text=None, extra_text=None,
                       *, fill=GOLD, text_color=NAVY,
                       top_size=14, bottom_size=10):
    """12-point star background + a separate text-box overlay.

    The star is purely decorative; the text lives in a normal
    rectangular text box layered on top so the text isn't constrained
    by the star's geometry (no risk of bleeding into the points).
    Reusable on every slide where MB = MC is being invoked, so the same
    visual pattern carries over.
    """
    left, top, width, height = int(left), int(top), int(width), int(height)

    # 1. Decorative star background (no text)
    shp = slide.shapes.add_shape(MSO_SHAPE.STAR_12_POINT, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = NAVY
    shp.line.width = Pt(1.0)
    shp.shadow.inherit = False
    # Soft drop shadow — added 2026-05-15 per user request, so the MB=MC
    # star reads as lifted off the slide like every other content shape.
    _add_drop_shadow(shp)
    # Suppress any auto-inserted text frame contents on the shape itself.
    shp.text_frame.text = ""

    # 2. Overlay text box, sized to the star's inscribed body
    inner_w = int(width * 0.65)
    inner_h = int(height * 0.55)
    inner_x = left + (width - inner_w) // 2
    inner_y = top + (height - inner_h) // 2

    box = slide.shapes.add_textbox(inner_x, inner_y, inner_w, inner_h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = top_text
    r1.font.name = 'Calibri'
    r1.font.size = Pt(top_size)
    r1.font.bold = True
    r1.font.color.rgb = text_color
    if bottom_text:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = bottom_text
        r2.font.name = 'Calibri'
        r2.font.size = Pt(bottom_size)
        r2.font.italic = True
        r2.font.bold = True
        r2.font.color.rgb = text_color
    if extra_text:
        p3 = tf.add_paragraph()
        p3.alignment = PP_ALIGN.CENTER
        r3 = p3.add_run()
        r3.text = extra_text
        r3.font.name = 'Calibri'
        r3.font.size = Pt(bottom_size)
        r3.font.italic = True
        r3.font.bold = True
        r3.font.color.rgb = text_color
    return shp


# --------------------------------------------------------------------------
# OMML (Office Math Markup Language) equation helper – gives formulas a
# proper TeX-style render with italic variables, stacked fractions, real
# subscripts/superscripts.  Uses Cambria Math (the standard PPT math font).
# --------------------------------------------------------------------------

M_NS = 'http://schemas.openxmlformats.org/officeDocument/2006/math'
A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
A14_NS = 'http://schemas.microsoft.com/office/drawing/2010/main'


def _omml_fill(color):
    """Build the inner ``<a:solidFill>`` clause for an OMML run's rPr.

    ``color`` may be an ``RGBColor`` instance (or a 3-tuple of ints).
    Returns an empty string when ``color`` is None so callers can splice
    the result into rPr unconditionally.
    """
    if color is None:
        return ''
    return (
        f'<a:solidFill><a:srgbClr val="'
        f'{color[0]:02X}{color[1]:02X}{color[2]:02X}'
        f'"/></a:solidFill>'
    )


def _omml_run(text, *, color=None):
    """OMML run for an italic variable (default math style).

    Inside an oMath, italic style is the math default for Latin letters;
    we leave m:rPr out entirely so the Cambria Math italic comes through.
    The a:rPr applies drawing-level font sizing/coloring.  Pass ``color``
    to tint the run (e.g., green ΔL / ΔQ in the slide-14 Convention box).
    """
    return (
        f'<m:r xmlns:m="{M_NS}">'
        f'<a:rPr xmlns:a="{A_NS}" lang="en-US" b="0" i="1">'
        f'{_omml_fill(color)}'
        f'<a:latin typeface="Cambria Math"/>'
        f'<a:ea typeface="Cambria Math"/>'
        f'</a:rPr>'
        f'<m:t>{text}</m:t>'
        f'</m:r>'
    )


def _omml_text(text, *, color=None):
    """Upright-style OMML run (for operators, numbers, acronyms).

    Force plain (upright) style via <m:rPr><m:sty m:val="p"/></m:rPr> – this
    is the documented way to disable the math-default italics for the
    enclosed run.  Pass ``color`` to tint the run.
    """
    return (
        f'<m:r xmlns:m="{M_NS}">'
        f'<m:rPr><m:sty m:val="p"/></m:rPr>'
        f'<a:rPr xmlns:a="{A_NS}" lang="en-US" b="0" i="0">'
        f'{_omml_fill(color)}'
        f'<a:latin typeface="Cambria Math"/>'
        f'</a:rPr>'
        f'<m:t xml:space="preserve">{text}</m:t>'
        f'</m:r>'
    )


def _omml_sub(base, sub):
    """OMML subscript: base with subscript expression."""
    return (
        f'<m:sSub xmlns:m="{M_NS}">'
        f'<m:sSubPr><m:ctrlPr>'
        f'<a:rPr xmlns:a="{A_NS}" lang="en-US" i="1">'
        f'<a:latin typeface="Cambria Math"/></a:rPr>'
        f'</m:ctrlPr></m:sSubPr>'
        f'<m:e>{base}</m:e>'
        f'<m:sub>{sub}</m:sub>'
        f'</m:sSub>'
    )


def _omml_frac(num, den):
    """OMML stacked fraction: num / den."""
    return (
        f'<m:f xmlns:m="{M_NS}">'
        f'<m:fPr><m:ctrlPr>'
        f'<a:rPr xmlns:a="{A_NS}" lang="en-US" i="1">'
        f'<a:latin typeface="Cambria Math"/></a:rPr>'
        f'</m:ctrlPr></m:fPr>'
        f'<m:num>{num}</m:num>'
        f'<m:den>{den}</m:den>'
        f'</m:f>'
    )


def _omml_sup(base, sup):
    """OMML superscript: base^sup (e.g. Q²)."""
    return (
        f'<m:sSup xmlns:m="{M_NS}">'
        f'<m:sSupPr><m:ctrlPr>'
        f'<a:rPr xmlns:a="{A_NS}" lang="en-US" i="1">'
        f'<a:latin typeface="Cambria Math"/></a:rPr>'
        f'</m:ctrlPr></m:sSupPr>'
        f'<m:e>{base}</m:e>'
        f'<m:sup>{sup}</m:sup>'
        f'</m:sSup>'
    )


def _add_dashed_gridlines(axis_el):
    """Dashed light-grey major gridlines on an axis (Cobb-Douglas style)."""
    for old in axis_el.findall(qn('c:majorGridlines')):
        axis_el.remove(old)
    gl = ET.Element(qn('c:majorGridlines'))
    sp = ET.SubElement(gl, qn('c:spPr'))
    ln = ET.SubElement(sp, qn('a:ln'))
    ln.set('w', '9525')
    ln.set('cap', 'flat'); ln.set('cmpd', 'sng'); ln.set('algn', 'ctr')
    fill = ET.SubElement(ln, qn('a:solidFill'))
    clr = ET.SubElement(fill, qn('a:srgbClr')); clr.set('val', 'C8CDD3')
    dash = ET.SubElement(ln, qn('a:prstDash')); dash.set('val', 'dash')
    axpos = axis_el.find(qn('c:axPos'))
    if axpos is not None:
        axpos.addnext(gl)


def _align_x_labels_with_ticks(value_axis):
    """Set <c:crossBetween val="midCat"/> on a value axis so the category
    labels and tick marks align (default for line charts in OOXML is
    "between", which leaves labels visually between adjacent ticks).
    """
    val_el = value_axis._element
    for old in val_el.findall(qn('c:crossBetween')):
        val_el.remove(old)
    cb = ET.Element(qn('c:crossBetween'))
    cb.set('val', 'midCat')
    # Schema position: c:crossBetween follows c:crosses(At) and precedes
    # c:majorUnit.  Insert before c:majorUnit if present; else append.
    mu_el = val_el.find(qn('c:majorUnit'))
    if mu_el is not None:
        mu_el.addprevious(cb)
    else:
        val_el.append(cb)


def _make_simple_line_chart(slide, x, y, w, h, categories, values, *,
                              line_color, x_title, y_title,
                              y_min=0, y_max=None, y_unit=None,
                              marker='circle'):
    """Single-series line+markers chart with dashed light-grey gridlines.

    Same visual conventions as slide 11 (Calibri navy labels, dashed C8CDD3
    major gridlines, marker size 7) but no legend/title.
    """
    cd = CategoryChartData()
    cd.categories = list(categories)
    cd.add_series("Y", values)
    # Drop-shadow rectangle behind the chart (graphicFrames can't host shadow).
    _add_graphicframe_shadow(slide, x, y, w, h)
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE,
        int(x), int(y), int(w), int(h), cd,
    )
    chart = chart_shape.chart
    chart.has_title = False
    chart.has_legend = False

    clr_hex = f'{line_color[0]:02X}{line_color[1]:02X}{line_color[2]:02X}'

    for series in chart.series:
        line = series.format.line
        line.color.rgb = line_color
        line.width = Pt(2.5)
        ser_xml = series._element
        for old in ser_xml.findall(qn('c:marker')):
            ser_xml.remove(old)
        m = ET.SubElement(ser_xml, qn('c:marker'))
        sym = ET.SubElement(m, qn('c:symbol')); sym.set('val', marker)
        sz_el = ET.SubElement(m, qn('c:size')); sz_el.set('val', '7')
        sp = ET.SubElement(m, qn('c:spPr'))
        fl = ET.SubElement(sp, qn('a:solidFill'))
        rg = ET.SubElement(fl, qn('a:srgbClr')); rg.set('val', clr_hex)
        ln = ET.SubElement(sp, qn('a:ln'))
        lf = ET.SubElement(ln, qn('a:solidFill'))
        lr = ET.SubElement(lf, qn('a:srgbClr')); lr.set('val', clr_hex)
        # disable smoothing (straight segments between points)
        for sm in ser_xml.findall(qn('c:smooth')):
            ser_xml.remove(sm)
        smooth = ET.SubElement(ser_xml, qn('c:smooth'))
        smooth.set('val', '0')

    # Axes – axis titles in BOLD ITALIC navy (per course CLAUDE.md);
    # tick labels in regular Calibri navy.
    cat = chart.category_axis
    cat.tick_labels.font.name = "Calibri"
    cat.tick_labels.font.size = Pt(10)
    cat.tick_labels.font.color.rgb = NAVY
    cat.has_title = True
    cat.axis_title.text_frame.text = x_title
    ar = cat.axis_title.text_frame.paragraphs[0].runs[0]
    ar.font.name = "Calibri"; ar.font.size = Pt(12)
    ar.font.bold = True; ar.font.italic = True
    ar.font.color.rgb = NAVY

    val = chart.value_axis
    val.tick_labels.font.name = "Calibri"
    val.tick_labels.font.size = Pt(10)
    val.tick_labels.font.color.rgb = NAVY
    val.has_title = True
    val.axis_title.text_frame.text = y_title
    ar = val.axis_title.text_frame.paragraphs[0].runs[0]
    ar.font.name = "Calibri"; ar.font.size = Pt(12)
    ar.font.bold = True; ar.font.italic = True
    ar.font.color.rgb = NAVY
    if y_max is not None:
        val.minimum_scale = y_min
        val.maximum_scale = y_max
    if y_unit is not None:
        val.major_unit = y_unit

    # Align X-axis category labels with tick marks (default OOXML places
    # them in the gaps between ticks).
    _align_x_labels_with_ticks(val)

    _add_dashed_gridlines(cat._element)
    _add_dashed_gridlines(val._element)
    return chart_shape


def _make_multi_line_chart(slide, x, y, w, h, categories, series, *,
                             x_title, y_title,
                             y_min=0, y_max=None, y_unit=None,
                             legend=True, legend_pos=('0.08', '0.10', '0.20', '0.20')):
    """Multi-series line+markers chart with the deck's standard styling.

    series: list of (name, values, color: RGBColor, marker: str) tuples.
    legend_pos: (x, y, w, h) in chart-fraction units (str) – top-left default.
    """
    _add_graphicframe_shadow(slide, x, y, w, h)
    cd = CategoryChartData()
    cd.categories = list(categories)
    for name, values, _color, _marker in series:
        cd.add_series(name, values)
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE,
        int(x), int(y), int(w), int(h), cd,
    )
    chart = chart_shape.chart
    chart.has_title = False
    chart.has_legend = legend

    # Per-series styling: line color + marker
    for idx, ser in enumerate(chart.series):
        name, values, color, marker = series[idx]
        line = ser.format.line
        line.color.rgb = color
        line.width = Pt(2.5)
        ser_xml = ser._element
        clr_hex = f'{color[0]:02X}{color[1]:02X}{color[2]:02X}'
        for old in ser_xml.findall(qn('c:marker')):
            ser_xml.remove(old)
        m = ET.SubElement(ser_xml, qn('c:marker'))
        sym = ET.SubElement(m, qn('c:symbol')); sym.set('val', marker)
        sz_el = ET.SubElement(m, qn('c:size')); sz_el.set('val', '7')
        sp = ET.SubElement(m, qn('c:spPr'))
        fl = ET.SubElement(sp, qn('a:solidFill'))
        rg = ET.SubElement(fl, qn('a:srgbClr')); rg.set('val', clr_hex)
        ln = ET.SubElement(sp, qn('a:ln'))
        lf = ET.SubElement(ln, qn('a:solidFill'))
        lr = ET.SubElement(lf, qn('a:srgbClr')); lr.set('val', clr_hex)
        for sm in ser_xml.findall(qn('c:smooth')):
            ser_xml.remove(sm)
        smooth = ET.SubElement(ser_xml, qn('c:smooth'))
        smooth.set('val', '0')

    # Axes – bold italic navy titles per course style
    cat = chart.category_axis
    cat.tick_labels.font.name = "Calibri"
    cat.tick_labels.font.size = Pt(10)
    cat.tick_labels.font.color.rgb = NAVY
    cat.has_title = True
    cat.axis_title.text_frame.text = x_title
    ar = cat.axis_title.text_frame.paragraphs[0].runs[0]
    ar.font.name = "Calibri"; ar.font.size = Pt(12)
    ar.font.bold = True; ar.font.italic = True; ar.font.color.rgb = NAVY

    val = chart.value_axis
    val.tick_labels.font.name = "Calibri"
    val.tick_labels.font.size = Pt(10)
    val.tick_labels.font.color.rgb = NAVY
    val.has_title = True
    val.axis_title.text_frame.text = y_title
    ar = val.axis_title.text_frame.paragraphs[0].runs[0]
    ar.font.name = "Calibri"; ar.font.size = Pt(12)
    ar.font.bold = True; ar.font.italic = True; ar.font.color.rgb = NAVY
    if y_max is not None:
        val.minimum_scale = y_min
        val.maximum_scale = y_max
    if y_unit is not None:
        val.major_unit = y_unit

    # Align X-axis category labels with tick marks.
    _align_x_labels_with_ticks(val)

    _add_dashed_gridlines(cat._element)
    _add_dashed_gridlines(val._element)

    # Legend top-left inside plot with white fill
    if legend:
        leg_el = chart.legend._element
        chart.legend.font.name = "Calibri"
        chart.legend.font.size = Pt(11)
        chart.legend.font.color.rgb = NAVY
        chart.legend.include_in_layout = False
        # Strip default legendPos / layout, replace
        for old in leg_el.findall(qn('c:layout')):
            leg_el.remove(old)
        for old in leg_el.findall(qn('c:legendPos')):
            leg_el.remove(old)
        pos = ET.SubElement(leg_el, qn('c:legendPos')); pos.set('val', 'tr')
        leg_el.remove(pos); leg_el.insert(0, pos)
        # manualLayout positions in chart-fraction units
        layout = ET.Element(qn('c:layout'))
        ml = ET.SubElement(layout, qn('c:manualLayout'))
        ET.SubElement(ml, qn('c:xMode')).set('val', 'edge')
        ET.SubElement(ml, qn('c:yMode')).set('val', 'edge')
        ET.SubElement(ml, qn('c:x')).set('val', legend_pos[0])
        ET.SubElement(ml, qn('c:y')).set('val', legend_pos[1])
        ET.SubElement(ml, qn('c:w')).set('val', legend_pos[2])
        ET.SubElement(ml, qn('c:h')).set('val', legend_pos[3])
        pos.addnext(layout)
        # White fill behind legend
        for old in leg_el.findall(qn('c:spPr')):
            leg_el.remove(old)
        leg_spPr = ET.Element(qn('c:spPr'))
        sf = ET.SubElement(leg_spPr, qn('a:solidFill'))
        clr = ET.SubElement(sf, qn('a:srgbClr')); clr.set('val', 'FFFFFF')
        ln = ET.SubElement(leg_spPr, qn('a:ln')); ln.set('w', '6350')
        lf = ET.SubElement(ln, qn('a:solidFill'))
        lc = ET.SubElement(lf, qn('a:srgbClr')); lc.set('val', '0B2B4E')
        layout.addnext(leg_spPr)
    return chart_shape


def _make_xy_line_chart(slide, x, y, w, h, *, series, x_title, y_title,
                          x_min=0, x_max=None, x_unit=None,
                          y_min=0, y_max=None, y_unit=None,
                          legend=False, legend_pos=None, smooth=False):
    """XY-scatter line+markers chart with the deck's standard styling.

    Use when you need data points to sit at arbitrary X-positions (e.g.,
    plotting MPL at the midpoint of each L-interval) while keeping tick
    marks at standard L-values.  Both axes are numeric value axes.

    series: list of ``(name, [(x, y), ...], color: RGBColor, marker: str)``.
    smooth: True for XY_SCATTER_SMOOTH (cubic spline through points),
            False for XY_SCATTER_LINES (straight segments).
    """
    _add_graphicframe_shadow(slide, x, y, w, h)
    cd = XyChartData()
    for name, points, _color, _marker in series:
        s = cd.add_series(name)
        for px, py in points:
            s.add_data_point(px, py)
    chart_type = (XL_CHART_TYPE.XY_SCATTER_SMOOTH if smooth
                  else XL_CHART_TYPE.XY_SCATTER_LINES)
    chart_shape = slide.shapes.add_chart(
        chart_type,
        int(x), int(y), int(w), int(h), cd,
    )
    chart = chart_shape.chart
    chart.has_title = False
    chart.has_legend = legend

    # Per-series styling
    for idx, ser in enumerate(chart.series):
        name, points, color, marker = series[idx]
        line = ser.format.line
        line.color.rgb = color
        line.width = Pt(2.5)
        ser_xml = ser._element
        clr_hex = f'{color[0]:02X}{color[1]:02X}{color[2]:02X}'
        for old in ser_xml.findall(qn('c:marker')):
            ser_xml.remove(old)
        m = ET.SubElement(ser_xml, qn('c:marker'))
        sym = ET.SubElement(m, qn('c:symbol')); sym.set('val', marker)
        sz_el = ET.SubElement(m, qn('c:size')); sz_el.set('val', '7')
        sp = ET.SubElement(m, qn('c:spPr'))
        fl = ET.SubElement(sp, qn('a:solidFill'))
        rg = ET.SubElement(fl, qn('a:srgbClr')); rg.set('val', clr_hex)
        ln = ET.SubElement(sp, qn('a:ln'))
        lf = ET.SubElement(ln, qn('a:solidFill'))
        lr = ET.SubElement(lf, qn('a:srgbClr')); lr.set('val', clr_hex)
        for sm in ser_xml.findall(qn('c:smooth')):
            ser_xml.remove(sm)
        smooth = ET.SubElement(ser_xml, qn('c:smooth'))
        smooth.set('val', '0')

    # Both axes are value axes in XY scatter; python-pptx returns the
    # X axis through .category_axis (wrapped as a ValueAxis since
    # catAx_lst is empty) and the Y axis through .value_axis.
    x_ax = chart.category_axis
    x_ax.tick_labels.font.name = "Calibri"
    x_ax.tick_labels.font.size = Pt(10)
    x_ax.tick_labels.font.color.rgb = NAVY
    if x_max is not None:
        x_ax.minimum_scale = x_min
        x_ax.maximum_scale = x_max
    if x_unit is not None:
        x_ax.major_unit = x_unit
    x_ax.has_title = True
    x_ax.axis_title.text_frame.text = x_title
    ar = x_ax.axis_title.text_frame.paragraphs[0].runs[0]
    ar.font.name = "Calibri"; ar.font.size = Pt(12)
    ar.font.bold = True; ar.font.italic = True; ar.font.color.rgb = NAVY

    y_ax = chart.value_axis
    y_ax.tick_labels.font.name = "Calibri"
    y_ax.tick_labels.font.size = Pt(10)
    y_ax.tick_labels.font.color.rgb = NAVY
    if y_max is not None:
        y_ax.minimum_scale = y_min
        y_ax.maximum_scale = y_max
    if y_unit is not None:
        y_ax.major_unit = y_unit
    y_ax.has_title = True
    y_ax.axis_title.text_frame.text = y_title
    ar = y_ax.axis_title.text_frame.paragraphs[0].runs[0]
    ar.font.name = "Calibri"; ar.font.size = Pt(12)
    ar.font.bold = True; ar.font.italic = True; ar.font.color.rgb = NAVY

    _add_dashed_gridlines(x_ax._element)
    _add_dashed_gridlines(y_ax._element)

    # Legend
    if legend and legend_pos is not None:
        leg_el = chart.legend._element
        chart.legend.font.name = "Calibri"
        chart.legend.font.size = Pt(11)
        chart.legend.font.color.rgb = NAVY
        chart.legend.include_in_layout = False
        for old in leg_el.findall(qn('c:layout')):
            leg_el.remove(old)
        for old in leg_el.findall(qn('c:legendPos')):
            leg_el.remove(old)
        pos = ET.SubElement(leg_el, qn('c:legendPos')); pos.set('val', 'tr')
        leg_el.remove(pos); leg_el.insert(0, pos)
        layout = ET.Element(qn('c:layout'))
        ml = ET.SubElement(layout, qn('c:manualLayout'))
        ET.SubElement(ml, qn('c:xMode')).set('val', 'edge')
        ET.SubElement(ml, qn('c:yMode')).set('val', 'edge')
        ET.SubElement(ml, qn('c:x')).set('val', legend_pos[0])
        ET.SubElement(ml, qn('c:y')).set('val', legend_pos[1])
        ET.SubElement(ml, qn('c:w')).set('val', legend_pos[2])
        ET.SubElement(ml, qn('c:h')).set('val', legend_pos[3])
        pos.addnext(layout)
        for old in leg_el.findall(qn('c:spPr')):
            leg_el.remove(old)
        leg_spPr = ET.Element(qn('c:spPr'))
        sf = ET.SubElement(leg_spPr, qn('a:solidFill'))
        clr = ET.SubElement(sf, qn('a:srgbClr')); clr.set('val', 'FFFFFF')
        ln = ET.SubElement(leg_spPr, qn('a:ln')); ln.set('w', '6350')
        lf = ET.SubElement(ln, qn('a:solidFill'))
        lc = ET.SubElement(lf, qn('a:srgbClr')); lc.set('val', '0B2B4E')
        layout.addnext(leg_spPr)

    return chart_shape


def _omml_acc_overline(symbol):
    """Inline OMML accent (overline / bar) on a math symbol.

    symbol: e.g. 'K' or 'L' – the variable to wear the bar.
    Returns an OMML fragment intended to be embedded inside an <m:oMath>.
    """
    return (
        '<m:acc>'
          '<m:accPr><m:chr m:val="̅"/></m:accPr>'
          '<m:e>'
            '<m:r>'
              '<a:rPr lang="en-US" b="0" i="1">'
                '<a:latin typeface="Cambria Math"/>'
              '</a:rPr>'
              f'<m:t>{symbol}</m:t>'
            '</m:r>'
          '</m:e>'
        '</m:acc>'
    )


def _add_mixed_textbox(slide, left, top, width, height, segments, *,
                        align=PP_ALIGN.LEFT, default_color=NAVY,
                        default_size=24,
                        margin_left=None, margin_right=None,
                        margin_top=None, margin_bottom=None,
                        space_before_pts=None):
    """Build a textbox whose paragraphs mix plain text runs and inline OMML.

    segments: list of (kind, content, opts) tuples, with kind ∈ {"text",
    "omml", "break"}.  "break" inserts a new paragraph.  Opts may set
    `size`, `bold`, `italic`, `color`, `font` per run.

    ``space_before_pts`` puts the same spcBef on every paragraph — the
    sparse-slide spacing rule (few lines on a slide breathe better with
    18 pt between them; Teaching CLAUDE.md, 2026-08-26).
    """
    box = slide.shapes.add_textbox(int(left), int(top),
                                     int(width), int(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05) if margin_left is None else margin_left
    tf.margin_right = Inches(0.05) if margin_right is None else margin_right
    tf.margin_top = Inches(0.0) if margin_top is None else margin_top
    tf.margin_bottom = Inches(0.0) if margin_bottom is None else margin_bottom

    align_attr = ''
    if align == PP_ALIGN.CENTER: align_attr = ' algn="ctr"'
    elif align == PP_ALIGN.RIGHT: align_attr = ' algn="r"'

    spc_xml = ('<a:spcBef><a:spcPts val="%d"/></a:spcBef>'
               % int(space_before_pts * 100)) if space_before_pts else ''

    def _start_para():
        if spc_xml:
            ppr = f'<a:pPr{align_attr}>{spc_xml}</a:pPr>'
        elif align_attr:
            ppr = f'<a:pPr{align_attr}/>'
        else:
            ppr = ''
        return [f'<a:p xmlns:a="{A_NS}" xmlns:m="{M_NS}" xmlns:a14="{A14_NS}">',
                ppr]

    paragraphs = [_start_para()]
    for kind, content, opts in segments:
        if kind == 'break':
            paragraphs[-1].append('<a:endParaRPr lang="en-US"/></a:p>')
            paragraphs.append(_start_para())
            continue
        size_pt = int(opts.get('size', default_size) * 100)
        color = opts.get('color', default_color)
        clr_hex = f'{color[0]:02X}{color[1]:02X}{color[2]:02X}'
        bold_attr = ' b="1"' if opts.get('bold') else ''
        italic_attr = ' i="1"' if opts.get('italic') else ''
        font = opts.get('font', 'Calibri')
        if kind == 'text':
            paragraphs[-1].append(
                f'<a:r><a:rPr lang="en-US" sz="{size_pt}"{bold_attr}{italic_attr}>'
                f'<a:solidFill><a:srgbClr val="{clr_hex}"/></a:solidFill>'
                f'<a:latin typeface="{font}"/>'
                f'</a:rPr><a:t>{content}</a:t></a:r>'
            )
        elif kind == 'omml':
            paragraphs[-1].append(
                f'<a14:m><m:oMath>{content}</m:oMath></a14:m>'
            )
    paragraphs[-1].append('<a:endParaRPr lang="en-US"/></a:p>')
    full_xml = ''.join(''.join(p) for p in paragraphs)

    txBody = tf._txBody
    for old in list(txBody.findall(qn('a:p'))):
        txBody.remove(old)
    # We have to parse each <a:p> separately so the namespaces resolve
    for p_str in full_xml.split('</a:p>'):
        if not p_str.strip(): continue
        p_xml = p_str + '</a:p>'
        new_p = ET.fromstring(p_xml)
        txBody.append(new_p)
        # Apply size+color to any OMML m:r elements inside this paragraph.
        # Color is set to ``default_color`` ONLY when no per-run solidFill
        # is already present — this lets callers tint individual OMML runs
        # via the optional ``color=`` argument on ``_omml_run`` / ``_omml_text``
        # (e.g., the green ΔL / ΔQ in the slide-14 Convention box) without
        # being silently overridden here.
        clr_hex = f'{default_color[0]:02X}{default_color[1]:02X}{default_color[2]:02X}'
        for r in new_p.iter(qn('m:r')):
            arPr = r.find(qn('a:rPr'))
            if arPr is None:
                arPr = ET.Element(qn('a:rPr'))
                arPr.set('lang', 'en-US')
                r.insert(0, arPr)
            arPr.set('sz', str(int(default_size * 100)))
            if arPr.find(qn('a:solidFill')) is None:
                # fill BEFORE a:latin (schema order) or the color is ignored
                sf = arPr.makeelement(qn('a:solidFill'), {})
                srgb = ET.SubElement(sf, qn('a:srgbClr'))
                srgb.set('val', clr_hex)
                arPr.insert(0, sf)
    return box


def _add_math_equation(slide, left, top, width, height, omml_content, *,
                       size_pt=32, color=NAVY, fill=None, line=None,
                       rounded=False, shadow=None, corner_pct=25000):
    """Place an OMML equation in a textbox on the slide.

    omml_content: a string built from _omml_* helpers (without the outer
    <m:oMathPara> wrapper).

    ``rounded=True`` gives the fill box rounded corners (``corner_pct`` in
    OOXML adj units, 25000 ≈ 25%); ``shadow=True`` adds a soft drop
    shadow.  Both default off.  IMPORTANT: the roundRect prstGeom must be
    inserted right after <a:xfrm> and BEFORE the fill, or PowerPoint
    silently drops the shape (see slide-54 fix).
    """
    left, top, width, height = int(left), int(top), int(width), int(height)
    # 2026-08-24 (Nico): a FILLED equation box is one of the deck's cream
    # cards, so it gets the soft shade unless the caller says otherwise.
    # Passing shadow=False still suppresses it.
    if shadow is None:
        shadow = fill is not None
    box = slide.shapes.add_textbox(left, top, width, height)

    if fill is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
    if line is not None:
        box.line.color.rgb = line
        box.line.width = Pt(0.75)

    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Replace the default <a:p> with one that hosts a math zone (<a14:m>)
    # containing the oMathPara.  The <a14:m> wrapper is REQUIRED by
    # PowerPoint to recognise OMML inside a textbox – without it PPT just
    # shows empty boxes.
    txBody = tf._txBody
    for p in list(txBody.findall(qn('a:p'))):
        txBody.remove(p)

    sz = int(size_pt * 100)
    clr_hex = '{:02X}{:02X}{:02X}'.format(color[0], color[1], color[2])

    # Use unique namespace prefixes; lxml will resolve them when parsing.
    p_xml = (
        f'<a:p xmlns:a="{A_NS}" xmlns:m="{M_NS}" xmlns:a14="{A14_NS}">'
        f'<a:pPr algn="ctr">'
        f'<a:defRPr sz="{sz}" b="0" i="1">'
        f'<a:solidFill><a:srgbClr val="{clr_hex}"/></a:solidFill>'
        f'<a:latin typeface="Cambria Math"/>'
        f'</a:defRPr>'
        f'</a:pPr>'
        f'<a14:m>'
        f'<m:oMathPara>'
        f'<m:oMathParaPr><m:jc m:val="centerGroup"/></m:oMathParaPr>'
        f'<m:oMath>{omml_content}</m:oMath>'
        f'</m:oMathPara>'
        f'</a14:m>'
        f'<a:endParaRPr lang="en-US" sz="{sz}"/>'
        f'</a:p>'
    )

    new_p = ET.fromstring(p_xml)
    txBody.append(new_p)

    # Set size and color on every OMML run's <a:rPr> so the text renders at
    # the right size.  (The defRPr above is the fallback.)
    for r in new_p.iter(qn('m:r')):
        arPr = r.find(qn('a:rPr'))
        if arPr is None:
            arPr = ET.Element(qn('a:rPr'))
            r.insert(0, arPr)
        arPr.set('sz', str(sz))
        if arPr.get('lang') is None:
            arPr.set('lang', 'en-US')
        # respect per-run colors set via _omml_run/_omml_text(color=...);
        # only runs WITHOUT a fill get the equation's default color
        # (was clobbering firm-color runs — Nico 2026-08-07)
        if not arPr.findall(qn('a:solidFill')):
            sf = ET.Element(qn('a:solidFill'))
            srgb = ET.SubElement(sf, qn('a:srgbClr'))
            srgb.set('val', clr_hex)
            arPr.insert(0, sf)

    # Optional rounded corners + drop shadow on the fill box.  prstGeom
    # MUST sit right after <a:xfrm> and before <a:solidFill>, else PPT
    # silently refuses to render the shape (slide-54 lesson).
    if rounded or shadow:
        spPr = box._element.find(qn('p:spPr'))
        if spPr is not None:
            if rounded:
                for old in spPr.findall(qn('a:prstGeom')):
                    spPr.remove(old)
                prstGeom = ET.Element(qn('a:prstGeom'))
                prstGeom.set('prst', 'roundRect')
                avLst = ET.SubElement(prstGeom, qn('a:avLst'))
                gd = ET.SubElement(avLst, qn('a:gd'))
                gd.set('name', 'adj'); gd.set('fmla', f'val {int(corner_pct)}')
                xfrm = spPr.find(qn('a:xfrm'))
                if xfrm is not None:
                    xfrm.addnext(prstGeom)
                else:
                    spPr.insert(0, prstGeom)
            if shadow:
                _add_drop_shadow(box)
    return box




# ==========================================================================
#  MODULE 2 — DEMAND ANALYSIS, IN-CLASS DECK  (everything below is M2)
#
#  Above this line: the reusable helper layer — do not edit here.
#
#  Pipeline (rerunnable, Module 7 pattern):
#    python _build_Module2InClass.py   -> full 76-slide deck, stubs for
#                                         polls + the pizza Excel slide
#    python _splice_media.py           -> verbatim OOXML splice of the 8
#                                         PollEv pairs + Excel slide from
#                                         "Module 2 - In Class with
#                                         Solutions.pptx"
#    python _animate.py all apply      -> fade builds per slide plans
# ==========================================================================

import uuid

CREAM = RGBColor(0xFD, 0xF6, 0xE6)
DIM = RGBColor(0xBF, 0xBF, 0xBF)           # 2026-08-24 (Nico):
                                           # outline items not
                                           # currently covered are
                                           # shaded (his own decks
                                           # use schemeClr bg1
                                           # lumMod 75% over white)
# How far a SHADED outline item's text drops so its single line centres
# against the gold circle (Nico, 2026-08-25).  Measured off the three
# final video decks, not guessed: every dimmed item moved down by
# exactly this much and the current topic did not move at all.
DIM_DROP = 85064                           # 0.093 in
# concept blue and a pale tint of it, for the boxed takeaway Nico asked
# for on slide 15 (2026-08-25) - deliberately NOT cream, so it reads as
# a different device from the convention callouts
CBLUE = RGBColor(0x00, 0x70, 0xC0)
PALE_BLUE = RGBColor(0xEA, 0xF3, 0xFC)
# dark red, for the slope callout Nico wanted in place of CT's
# gold on the point-elasticity slide (2026-08-25)
DARK_RED = RGBColor(0xA5, 0x1C, 0x1C)

# 2026-08-28 (Nico): the top-bar tag names the CURRENT AGENDA ITEM, so
# every content slide reads "Module 2 · <outline item title>" taken
# verbatim from M2_OUTLINE.  This SUPERSEDES the old three-level
# "Module 2 · Part N · Section" tag, whose middle level had drifted away
# from the agenda wording (Teaching CLAUDE.md, "Canvas, palette, and
# chrome").  Own-price / income / cross-price are sub-sections of
# outline item 2 ("Elasticities"), not agenda items of their own, so
# they all resolve to the same tag - the constants stay separate only so
# a future split of item 2 needs one edit here rather than 40 call sites.
# The module front matter (logistics, recap, course roadmap) keeps its
# own two-level tag; the outline slides read "Module 2 · Agenda".
TAG_LOGISTICS = "Module 2 · Logistics"
TAG_RECAP     = "Module 2 · Recap"
TAG_ROADMAP   = "Module 2 · Course Roadmap"
TAG_OUTLINE   = "Module 2 · Agenda"
TAG_LAW       = "Module 2 · The Law of Demand"
TAG_ELAST     = "Module 2 · Elasticities"
TAG_OWN       = "Module 2 · Elasticities"
TAG_INCOME    = "Module 2 · Elasticities"
TAG_CROSS     = "Module 2 · Elasticities"
# 2026-08-28: the cheat sheet is still an item-2 slide; the two
# post-work outline slides are agenda slides and take TAG_OUTLINE
TAG_WRAP      = "Module 2 · Elasticities"


# --------------------------------------------------------------------------
# M7 override: footer page number is a LIVE slide-number field
# (Teaching CLAUDE.md: "Footer page numbers are LIVE slide-number fields").
# Defined AFTER the helpers so every helper call resolves to this version.
# --------------------------------------------------------------------------

def _add_slidenum_field(slide, left, top, width, height, page_num, *,
                        size=12, color=GRAY):
    tb = slide.shapes.add_textbox(int(left), int(top), int(width), int(height))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    p_el = p._p
    guid = str(uuid.uuid5(uuid.NAMESPACE_DNS, f"m2ic-slidenum-{page_num}")).upper()
    fld = p_el.makeelement(qn('a:fld'), {'id': '{%s}' % guid, 'type': 'slidenum'})
    rPr = p_el.makeelement(qn('a:rPr'),
                           {'lang': 'en-US', 'sz': str(size * 100), 'dirty': '0'})
    fill = p_el.makeelement(qn('a:solidFill'), {})
    clr = p_el.makeelement(qn('a:srgbClr'), {'val': str(color)})
    fill.append(clr)
    rPr.append(fill)
    latin = p_el.makeelement(qn('a:latin'), {'typeface': 'Calibri'})
    rPr.append(latin)
    t = p_el.makeelement(qn('a:t'), {})
    t.text = str(page_num)
    fld.append(rPr)
    fld.append(t)
    p_el.append(fld)
    return tb


def _draw_footer(slide, footer_text, page_num):  # noqa: F811 — M7 override
    _add_rect(slide, 0, Inches(7.15), SLIDE_W, Inches(0.02), RULE)
    _add_rect(slide, MARGIN, Inches(7.135), GOLD_W, Inches(0.05), GOLD)
    _add_text(slide, MARGIN, Inches(7.20), Inches(11), Inches(0.32),
              footer_text, size=12, color=GRAY)
    _add_slidenum_field(slide, Inches(12.55), Inches(7.20), Inches(0.55),
                        Inches(0.32), page_num)


# --------------------------------------------------------------------------
# M7 image loader — source images live in _source_images/ under their
# original media names (imageNN.ext), extracted 2026-07-29.
# --------------------------------------------------------------------------

def _add_media_image(slide, fname, *, left, top, width=None, height=None,
                     rounded=True, shadow=True, corner_pct=8,
                     transparency=None):
    """Place a source-deck image by media filename. Logos/screenshots:
    rounded=False, shadow=False (flat exception)."""
    path = SRC_IMG_DIR / fname
    kwargs = {"left": int(left), "top": int(top)}
    if width is not None:
        kwargs["width"] = int(width)
    if height is not None:
        kwargs["height"] = int(height)
    pic = slide.shapes.add_picture(str(path), **kwargs)
    if transparency:
        # 2026-08-25: wash a full-bleed background photo out so dark
        # text stays readable on top (CT's lottery slide)
        blip = pic._element.find('.//' + qn('a:blip'))
        amf = ET.SubElement(blip, qn('a:alphaModFix'))
        amf.set('amt', str(int((100 - transparency) * 1000)))
    if rounded:
        _apply_picture_style(pic, corner_pct=corner_pct)
    elif shadow:
        _add_drop_shadow(pic)
    return pic


# --------------------------------------------------------------------------
# Native styled table (navy header, white/cream body, thin borders) on a
# shadowed white backing card.  NOTE: backing + table must be GROUPED in
# phase 3 (grpSp surgery) so the shade travels with the table.
# --------------------------------------------------------------------------

def _set_cell_borders(cell, *, color=RULE, weight_pt=0.75):
    tcPr = cell._tc.get_or_add_tcPr()
    for tag in ('a:lnL', 'a:lnR', 'a:lnT', 'a:lnB'):
        for old in tcPr.findall(qn(tag)):
            tcPr.remove(old)
    # schema order: lnL, lnR, lnT, lnB come first in tcPr
    for tag in reversed(('a:lnL', 'a:lnR', 'a:lnT', 'a:lnB')):
        ln = tcPr.makeelement(qn(tag), {'w': str(int(weight_pt * 12700)),
                                        'cap': 'flat'})
        fill = ln.makeelement(qn('a:solidFill'), {})
        c = fill.makeelement(qn('a:srgbClr'), {'val': str(color)})
        fill.append(c)
        ln.append(fill)
        tcPr.insert(0, ln)


def _add_styled_table(slide, left, top, width, height, rows_data, *,
                      col_widths=None, row_heights=None, font_size=18,
                      header_size=None, backing_pad=Inches(0.15),
                      first_col_bold=True, first_col_align_left=True,
                      cell_fills=None, cell_text_colors=None,
                      margin_v=None, backing=True,
                      backing_rounded=False):
    """rows_data: list of rows (row 0 = navy header). cell_fills /
    cell_text_colors: optional {(r, c): RGBColor} overrides.

    ``backing=False`` skips the shadowed white card — for a table that
    is one of several stacked pieces sharing ONE backing drawn by the
    caller (slide 52, where the three sections animate separately;
    three stacked cards would cast shadows across each other).
    """
    header_size = header_size or font_size
    left, top, width, height = int(left), int(top), int(width), int(height)
    if backing:
        _add_graphicframe_shadow(slide, left - int(backing_pad),
                                 top - int(backing_pad),
                                 width + 2 * int(backing_pad),
                                 height + 2 * int(backing_pad),
                                 rounded=backing_rounded)
    n_rows, n_cols = len(rows_data), len(rows_data[0])
    gf = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = gf.table
    tblPr = tbl._tbl.find(qn('a:tblPr'))
    if tblPr is not None:                      # strip default banding style
        tblPr.set('firstRow', '0')
        tblPr.set('bandRow', '0')
        for child in list(tblPr):
            tblPr.remove(child)
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = int(w)
    if row_heights:
        for i, h in enumerate(row_heights):
            tbl.rows[i].height = int(h)
    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.margin_left = Inches(0.10)
            cell.margin_right = Inches(0.10)
            # a long table can be squeezed by lowering margin_v; row
            # height in PowerPoint is line height + these two margins
            _mv = Inches(0.04) if margin_v is None else int(margin_v)
            cell.margin_top = _mv
            cell.margin_bottom = _mv
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            hdr = (r == 0)
            if cell_fills and (r, c) in cell_fills:
                cell.fill.solid()
                cell.fill.fore_color.rgb = cell_fills[(r, c)]
            elif hdr:
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else CREAM
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            if first_col_align_left and c == 0 and not hdr:
                p.alignment = PP_ALIGN.LEFT
            else:
                p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(val)
            run.font.name = "Calibri"
            run.font.size = Pt(header_size if hdr else font_size)
            run.font.bold = hdr or (first_col_bold and c == 0)
            if cell_text_colors and (r, c) in cell_text_colors:
                run.font.color.rgb = cell_text_colors[(r, c)]
            else:
                run.font.color.rgb = WHITE if hdr else NAVY
            _set_cell_borders(cell)
    return gf


# --------------------------------------------------------------------------
# M7 title slide, outline/divider, and stub
# --------------------------------------------------------------------------

def _inject_handoff_group(slide, fname, id_base=9500):
    """Append a hand-authored <p:grpSp> (saved verbatim from a canonical
    slide into _handoff_*.xml) into this slide's shape tree. Shape ids are
    re-based to avoid collisions. Used for Nico's scaled group cards
    (group scaling also scales text, which a rebuild can't replicate)."""
    xml = (OUT_DIR / fname).read_text(encoding='utf-8')
    el = ET.fromstring(xml)
    for i, nv in enumerate(el.iter(qn('p:cNvPr'))):
        nv.set('id', str(id_base + i))
    slide.shapes._spTree.append(el)
    return el



# --------------------------------------------------------------------------
# M2 footer + module-outline data
# --------------------------------------------------------------------------

FOOTER_TEXT = "Management 405  ·  Module 2  ·  Demand Analysis"

# (label, title, description, is_sub_item).  2026-08-24 (Nico): items 4
# and 5 became sub-items 3a and 3b of "Demand and revenue", and demand
# estimation moved up to 4 — the structure of the original deck, where
# those two sit at outline level 1 under "Demand and revenue".  The LIST
# INDICES are unchanged, so every highlight_idx / highlight_set call site
# keeps working.
M2_OUTLINE = [
    ("1", "The law of demand",
     "When the price falls, quantity demanded rises – holding everything "
     "else constant", False),
    ("2", "Elasticities",
     "How strongly demand responds to price, income, and the prices of "
     "other goods", False),
    ("3", "Demand and revenue",
     "How the demand curve translates into total revenue", False),
    ("3a", "Elasticity and revenue",
     "When a price increase raises revenue – and when it backfires", True),
    ("3b", "Marginal revenue",
     "The extra revenue from selling one more unit", True),
    ("4", "Demand estimation",
     "Measuring demand from data – market experiments and regression",
     False),
]

COURSE_PARTS = [
    "1.  Basic Concepts and Economic Principles",
    "2.  Value and Demand",
    "3.  Supply and Cost",
    "4.  Markets, Pricing and Strategy",
]

STUB_POLL = "PollEverywhere slide — spliced verbatim by _splice_media.py"
STUB_EXCEL = ("Live Excel-embed slide (class pizza demand) — spliced "
              "verbatim by _splice_media.py")


def _add_pollbreak_badge(slide):
    """Nico's hand-tuned Poll Break badge (2026-08-15): smaller gold
    parallelogram + navy label, grouped, bottom-right in FRONT of the
    footer. Injected verbatim from _handoff_pollbreak.xml (BUILD INPUT,
    never delete) — call AFTER _draw_footer."""
    _inject_handoff_group(slide, "_handoff_pollbreak.xml", id_base=9600)


# --------------------------------------------------------------------------
# M2 title slide
# --------------------------------------------------------------------------

def slide_01_title(prs):
    slide = _blank_slide(prs)
    _add_text(slide, 0, Inches(2.10), SLIDE_W, Inches(1.1),
              "Demand Analysis",
              size=60, bold=True, color=NAVY, font="Calibri",
              align=PP_ALIGN.CENTER)
    _add_text(slide, 0, Inches(3.25), SLIDE_W, Inches(0.75),
              "Module 2",
              size=40, bold=True, color=GOLD, font="Calibri",
              align=PP_ALIGN.CENTER)
    _add_rect(slide, int((SLIDE_W - Inches(4.0)) / 2), Inches(4.28),
              Inches(4.0), 54864, GOLD)
    _add_text(slide, 0, Inches(4.62), SLIDE_W, Inches(0.55),
              "Management 405", size=26, bold=True, color=GRAY,
              font="Calibri", align=PP_ALIGN.CENTER)
    _add_text(slide, 0, Inches(5.32), SLIDE_W, Inches(0.5),
              "Prof. Nico Voigtländer  ·  UCLA Anderson",
              size=22, color=GRAY, font="Calibri", align=PP_ALIGN.CENTER)
    _add_rect(slide, 0, Inches(7.15), SLIDE_W, Inches(0.02), RULE)
    _add_rect(slide, MARGIN, Inches(7.135), GOLD_W, Inches(0.05), GOLD)
    return slide


# --------------------------------------------------------------------------
# M2 module outline — one maker, two modes.
#   descriptions=True : the single descriptive overview (slide 7) — each of
#     the 6 items gets its title plus a one-line plain-language description.
#   highlight_idx / highlight_set : section-agenda mode — titles only, the
#     current item(s) navy on a cream band, the rest faded (CT style).
# --------------------------------------------------------------------------

# 2026-08-25 (Nico): every outline item says where it is taught.
# Topics 1 and 2 are done in class; 3, 3a, 3b and 4 in the videos.
#
# The pills name the video a student actually has to watch.  Topic 3
# spans TWO videos, so it says so; its sub-topics each name the single
# video they live in and get a narrower pill, right-aligned under the
# parent's, so they read as parts of it (Nico, 2026-08-25).
COVERAGE_LABEL = {
    0: "In class", 1: "In class",
    2: "Videos 1+2", 3: "Video 1", 4: "Video 2", 5: "Video 3",
}
IN_CLASS_ITEMS = {0, 1}


def make_m2_outline(prs, page_num, *, section_tag=TAG_OUTLINE,
                    title="Outline of Module 2", descriptions=False,
                    highlight_idx=None, highlight_set=None):
    """Module outline in CT's exact format (gold 0.58" circle at x=1.15,
    25 pt bold navy number + title, 22 pt gray description). Every item
    RESERVES the description row (uniform pitch, so item positions are
    identical on every agenda slide); the description text shows only
    for the current topic(s), or for all items when descriptions=True.
    On section agendas the current topic additionally gets a cream
    rounded band with gold border (no band on the descriptive
    overview/summary). Title y-positions land on slide._m2_item_ys."""
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, section_tag)
    _draw_action_title(slide, title)

    hi = set()
    if highlight_idx is not None:
        hi.add(highlight_idx)
    if highlight_set:
        hi.update(highlight_set)
    if descriptions:
        hi = set(range(len(M2_OUTLINE)))

    title_h = Inches(0.42)
    desc_h = Inches(0.38)
    gap = Inches(0.11)
    pitch = title_h + desc_h + gap
    total = pitch * len(M2_OUTLINE) - gap
    top = Inches(1.60)
    bottom = Inches(7.02)
    y = int(top + max(0, (bottom - top - total) // 2))

    ys = []
    for i, (label, item, desc, sub) in enumerate(M2_OUTLINE):
        ys.append(y)
        # sub-items (3a / 3b) sit indented under their parent, on a
        # smaller circle and one type size down
        circ_x = Inches(1.62) if sub else Inches(1.15)
        circ_d = Inches(0.46) if sub else Inches(0.58)
        circ_dy = Inches(0.08) if sub else Inches(0.02)
        text_x = Inches(2.28) if sub else Inches(2.05)
        num_pt = 17 if sub else 25
        item_pt = 22 if sub else 25
        band_x = Inches(1.37) if sub else Inches(0.90)
        band_w = Inches(11.68) if sub else Inches(12.15)
        # 2026-08-24 (Nico): on a section agenda the items that are NOT
        # currently covered are shaded; the descriptive overview, which
        # lights every item, keeps them all navy.  The gold circle fill
        # stays gold on every item either way.
        lit = descriptions or i in hi
        ink = NAVY if lit else DIM
        if not descriptions and i in hi:
            band = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, int(band_x),
                int(y - Inches(0.06)), int(band_w),
                int(title_h + desc_h + Inches(0.10)))
            band.adjustments[0] = 0.35
            band.fill.solid()
            band.fill.fore_color.rgb = CREAM
            band.line.color.rgb = GOLD
            band.line.width = Pt(1.0)
            band.shadow.inherit = False
            _add_drop_shadow(band)
        circ = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, int(circ_x), int(y + circ_dy),
            int(circ_d), int(circ_d))
        circ.fill.solid()
        circ.fill.fore_color.rgb = GOLD
        circ.line.fill.background()
        circ.shadow.inherit = False
        tf = circ.text_frame
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = label
        run.font.size = Pt(num_pt)
        run.font.bold = True
        run.font.color.rgb = ink
        run.font.name = "Calibri"
        # 2026-08-24 (Nico): agenda item titles follow the same title
        # case as slide titles.  The one-line description underneath is
        # a sentence, so it is left alone.
        rows = [([(_title_case(item[0].upper() + item[1:]),
                   {'bold': True, 'size': item_pt, 'color': ink})], 0,
                 {'bullet_style': 'none', 'space_before_pts': 0})]
        if i in hi:
            rows.append(([(desc, {'size': 22, 'color': GRAY})], 0,
                         {'bullet_style': 'none', 'space_before_pts': 0}))
        # 2026-08-25 (Nico, measured off the three final video decks):
        # an item with no description shown renders its single line at
        # the TOP of the reserved two-row box, which sits high against
        # the gold circle.  He nudged every SHADED item down by exactly
        # 85064 EMU (0.093 in) to centre it; the current topic, which
        # fills its box with title + description, does NOT move.
        _add_hierarchical_bullets(
            slide, text_x, y if lit else int(y + DIM_DROP), Inches(11.0),
            title_h + desc_h, rows, size=item_pt, line_spacing_pts=0)
        # where this topic is taught.  On a section agenda the pill
        # dims with the rest of the row - only the current topic keeps
        # its colour (Nico, 2026-08-25).
        in_class = i in IN_CLASS_ITEMS
        if lit:
            pill_fill = NAVY if in_class else GOLD
            pill_ink = WHITE if in_class else NAVY
        else:
            pill_fill, pill_ink = DIM, WHITE
        # a sub-item's pill is narrower and hangs off the same right
        # edge, so it reads as part of its parent's
        pill_w = Inches(1.14) if sub else Inches(1.55)
        pill_x = Inches(12.85) - pill_w
        _add_rounded_filled_box(
            slide, int(pill_x), int(y + Inches(0.02)),
            int(pill_w), Inches(0.36), COVERAGE_LABEL[i],
            fill=pill_fill, text_color=pill_ink,
            size=13, corner_pct=0.30, shadow=lit)
        y = int(y + pitch)

    slide._m2_item_ys = ys
    _draw_footer(slide, FOOTER_TEXT, page_num)
    return slide


# --------------------------------------------------------------------------
# Stub (positional placeholder for spliced poll / Excel slides)
# --------------------------------------------------------------------------

def make_stub(prs, page_num, section_tag, title, note, *, hidden=False):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, section_tag)
    _draw_action_title(slide, title)
    _add_text(slide, MARGIN, Inches(3.4), RULE_W, Inches(0.8),
              "[ %s ]" % note, size=20, italic=True, color=GRAY,
              font="Calibri", align=PP_ALIGN.CENTER)
    _draw_footer(slide, FOOTER_TEXT, page_num)
    if hidden:
        slide._element.set('show', '0')
    return slide


def make_todo(prs, page_num, section_tag, title):
    """Temporary placeholder for a slide not yet scripted — must be gone
    before the deck ships."""
    return make_stub(prs, page_num, section_tag, title,
                     "TODO — not yet scripted")




# ==========================================================================
#  BATCH B — front matter + Part 1: The Law of Demand (slides 2–24 core)
# ==========================================================================

RED = RGBColor(0xC0, 0x00, 0x00)          # deck red (worked-solution accents)
RED_FF = RGBColor(0xFF, 0x00, 0x00)       # original slide-9 c.p. red
GREEN = RGBColor(0x1B, 0x5E, 0x20)        # market/aggregate-demand green


class SimpleFig:
    """Logical→slide coordinate transform for shape-built charts.
    All geometry in float inches; returns Inches() EMU ints."""

    def __init__(self, left_in, bottom_in, w_in, h_in, xmax, ymax):
        self.l, self.b, self.w, self.h = left_in, bottom_in, w_in, h_in
        self.xmax, self.ymax = xmax, ymax

    def x(self, xv):
        return Inches(self.l + self.w * xv / self.xmax)

    def y(self, yv):
        return Inches(self.b - self.h * yv / self.ymax)


def _fig_axes(slide, fig, *, weight_pt=2.0):
    _add_arrow(slide, (fig.x(0), fig.y(0)),
               (fig.x(0), Inches(fig.b - fig.h - 0.18)),
               color=NAVY, weight_pt=weight_pt, head=True)
    _add_arrow(slide, (fig.x(0), fig.y(0)),
               (Inches(fig.l + fig.w + 0.18), fig.y(0)),
               color=NAVY, weight_pt=weight_pt, head=True)


def _fig_ytick(slide, fig, val, label, *, color=NAVY, size=16, bold=False):
    return _add_text(slide, Inches(fig.l - 1.07), fig.y(val) - Inches(0.14),
                     Inches(0.95), Inches(0.3), label, size=size,
                     bold=bold, color=color, font="Calibri",
                     align=PP_ALIGN.RIGHT)


def _fig_xtick(slide, fig, val, label, *, color=NAVY, size=16, bold=False):
    return _add_text(slide, fig.x(val) - Inches(0.5), Inches(fig.b + 0.06),
                     Inches(1.0), Inches(0.3), label, size=size,
                     bold=bold, color=color, font="Calibri",
                     align=PP_ALIGN.CENTER)


def _fig_point(slide, fig, xv, yv, *, r_in=0.055, fill=NAVY, line=None):
    d = Inches(2 * r_in)
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 fig.x(xv) - Inches(r_in),
                                 fig.y(yv) - Inches(r_in), d, d)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1.25)
    shp.shadow.inherit = False
    return shp


def _add_cubic_curve(slide, p0, c1, c2, p1, *, color=NAVY, weight_pt=2.5,
                     dash=None):
    """Editable cubic-Bézier freeform (custGeom, one segment, tight bbox).
    Points are (x, y) EMU/Inches tuples in slide coordinates."""
    pts = [(int(x), int(y)) for x, y in (p0, c1, c2, p1)]
    xs = [p[0] for p in pts]
    ys = [p[1] for p in pts]
    x0, y0 = min(xs), min(ys)
    w = max(max(xs) - x0, 1)
    h = max(max(ys) - y0, 1)
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x0, y0, w, h)
    shp.fill.background()
    shp.shadow.inherit = False
    sppr = shp._element.spPr
    prst = sppr.find(qn('a:prstGeom'))
    rel = [(p[0] - x0, p[1] - y0) for p in pts]
    cust = ET.fromstring(
        '<a:custGeom xmlns:a="http://schemas.openxmlformats.org/'
        'drawingml/2006/main">'
        '<a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/>'
        '<a:rect l="0" t="0" r="%d" b="%d"/>'
        '<a:pathLst><a:path w="%d" h="%d">'
        '<a:moveTo><a:pt x="%d" y="%d"/></a:moveTo>'
        '<a:cubicBezTo><a:pt x="%d" y="%d"/><a:pt x="%d" y="%d"/>'
        '<a:pt x="%d" y="%d"/></a:cubicBezTo>'
        '</a:path></a:pathLst></a:custGeom>'
        % (w, h, w, h, rel[0][0], rel[0][1], rel[1][0], rel[1][1],
           rel[2][0], rel[2][1], rel[3][0], rel[3][1]))
    sppr.replace(prst, cust)
    ln = sppr.find(qn('a:ln'))
    if ln is None:
        ln = ET.SubElement(sppr, qn('a:ln'))
    ln.set('w', str(int(weight_pt * 12700)))
    ln.set('cap', 'rnd')
    for child in list(ln):
        ln.remove(child)
    fill_el = ET.SubElement(ln, qn('a:solidFill'))
    clr = ET.SubElement(fill_el, qn('a:srgbClr'))
    clr.set('val', str(color))
    if dash:
        dash_el = ET.SubElement(ln, qn('a:prstDash'))
        dash_el.set('val', dash)
    ET.SubElement(ln, qn('a:round'))
    # schema order inside spPr is geometry -> fill -> ln -> effectLst.
    # shadow.inherit = False appends an empty <a:effectLst/> BEFORE the
    # ln we just built, and PowerPoint then silently drops the whole ln
    # and falls back to the theme line (thin, light blue).  Move any
    # effectLst behind the ln.  Found 2026-08-24 on the video deck's TR
    # parabola, which had been rendering unstyled.
    for eff in sppr.findall(qn('a:effectLst')):
        sppr.remove(eff)
        sppr.append(eff)
    return shp


# --------------------------------------------------------------------------
# Slide 2 — Logistics (placeholder dates per 2026-08-14 decision)
# --------------------------------------------------------------------------

def slide_02_logistics(prs):
    slide = make_content_bulleted(
        prs, 2, TAG_LOGISTICS, "Logistics",
        [
            ("Problem Set 1 due on [DATE] at 11:59pm", 0),
            ("Solutions: See email from TA", 1),
            ("Math Prerequisites: Take the short math test on BruinLearn", 0),
            ("Individual Midterm Assignment: [DATE] – [DATE]", 0),
            ("Flexible 3.5h time window", 1),
            ([("Actual time needed: Approx. 3hrs ", {}),
              ("(+30min for scanning and uploading)", {'size': 20})], 1, {}),
        ],
        size=26, sub_size=24)
    return slide


# --------------------------------------------------------------------------
# Slide 3 — Recap of Module 1
# --------------------------------------------------------------------------

def slide_03_recap(prs):
    return make_content_bulleted(
        prs, 3, TAG_RECAP, "Recap of Module 1",
        [
            ("Supply/demand framework for analyzing markets", 0),
            ("Key economic principles for decision-making", 0),
            ("Economic costs include opportunity costs", 1),
            ("Ignore sunk costs", 1),
            ("Use cost-benefit and marginal analysis", 1),
        ],
        size=28, sub_size=24)


# --------------------------------------------------------------------------
# Slide 6 — Agenda for the Class (course roadmap, "we are here" on Part 2)
# --------------------------------------------------------------------------

def slide_06_roadmap(prs):
    """Course roadmap in the Module-3 standard format (diamond layout,
    hand-approved 2026-08-15): module 2 highlighted navy, gold up-arrow
    "we are here" beneath it, faded boxes/connectors elsewhere."""

    def draw(slide):
        box_h = Inches(0.85)
        narrow_w = Inches(4.6)
        wide_w = Inches(8.6)
        gap = Inches(0.3)
        slide_mid = SLIDE_W // 2

        top_x = slide_mid - wide_w // 2
        top_y = Inches(2.0)
        _add_rounded_filled_box(
            slide, top_x, top_y, wide_w, box_h,
            "1. Basic Principles and Economic Way of Thinking",
            fill=FADED, text_color=WHITE, size=24, bold=True)

        row2_y = Inches(3.65)
        left_x = slide_mid - gap // 2 - narrow_w
        right_x = slide_mid + gap // 2
        _add_rounded_filled_box(slide, left_x, row2_y, narrow_w, box_h,
                                "2. Value and Demand",
                                fill=NAVY, text_color=WHITE, size=26,
                                bold=True)
        _add_rounded_filled_box(slide, right_x, row2_y, narrow_w, box_h,
                                "3. Supply and Cost",
                                fill=FADED, text_color=WHITE, size=26,
                                bold=True)

        bot_x = slide_mid - wide_w // 2
        bot_y = Inches(5.5)
        _add_rounded_filled_box(slide, bot_x, bot_y, wide_w, box_h,
                                "4. Markets, Pricing, and Strategy",
                                fill=FADED, text_color=WHITE, size=24,
                                bold=True)

        top_bottom_y = top_y + box_h
        _add_arrow(slide, (top_x + wide_w // 2, top_bottom_y),
                   (left_x + narrow_w // 2, row2_y),
                   color=NAVY, weight_pt=3.5, head=True)
        _add_arrow(slide, (top_x + wide_w // 2, top_bottom_y),
                   (right_x + narrow_w // 2, row2_y),
                   color=FADED, weight_pt=3.0, head=True)
        row2_bottom_y = row2_y + box_h
        _add_arrow(slide, (left_x + narrow_w // 2, row2_bottom_y),
                   (bot_x + wide_w // 2, bot_y),
                   color=FADED, weight_pt=3.0, head=True)
        _add_arrow(slide, (right_x + narrow_w // 2, row2_bottom_y),
                   (bot_x + wide_w // 2, bot_y),
                   color=FADED, weight_pt=3.0, head=True)

        # gold "we are here" arrow under the CURRENT box (module 2, left)
        left_shift = Inches(0.55)
        arrow_w = Inches(0.55)
        arrow_h = Inches(0.55)
        arrow_left = left_x + (narrow_w - arrow_w) // 2 - left_shift
        arrow_top = row2_y + box_h + Inches(0.05)
        _add_arrow_shape(slide, arrow_left, arrow_top, arrow_w, arrow_h,
                         direction="up", fill=GOLD)
        _add_text(slide, left_x - left_shift,
                  arrow_top + arrow_h + Inches(0.02),
                  narrow_w, Inches(0.32),
                  "we are here", size=16, italic=True, bold=True,
                  color=GOLD, font="Calibri", align=PP_ALIGN.CENTER)

    return make_diagram_slide(prs, 6, TAG_ROADMAP, "Agenda for the Class",
                              draw)


# --------------------------------------------------------------------------
# Slide 9 — The Law of Demand (+ native D-curve graphic, CT item 3)
# --------------------------------------------------------------------------

def slide_09_law_of_demand(prs, page_num=9):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_LAW)
    _draw_action_title(slide, "The Law of Demand")
    # Nico's hand restructure (2026-08-15): unbulleted section headers,
    # with boxes around the Assumption and Law-of-Demand sections
    _add_outlined_box(slide, Inches(0.30), Inches(1.62), Inches(8.30),
                      Inches(1.86), "", line=NAVY, fill=WHITE,
                      line_w=1.25, rounded=True, shadow=True,
                      corner_pct=0.08)
    _add_hierarchical_bullets(
        slide, Inches(0.55), Inches(1.76), Inches(7.85), Inches(1.52),
        [
            ("Crucial assumption:", 0,
             {'bold': True, 'bullet_style': 'none'}),
            ([("Holding everything else constant ", {}),
              ("(income, weather, product attributes, perception of "
               "product quality …)", {'size': 20})], 0, {}),
            ([("In econ parlance: “ceteris paribus” (", {}),
              ("c.p.", {'bold': True, 'color': RED_FF}),
              (")", {})], 0, {}),
        ],
        size=22, sub_size=20, line_spacing_pts=9)
    _add_outlined_box(slide, Inches(0.30), Inches(3.66), Inches(8.30),
                      Inches(1.88), "", line=NAVY, fill=WHITE,
                      line_w=1.25, rounded=True, shadow=True,
                      corner_pct=0.08)
    _add_hierarchical_bullets(
        slide, Inches(0.55), Inches(3.80), Inches(7.85), Inches(1.66),
        [
            ("The Law of Demand says:", 0,
             {'bold': True, 'bullet_style': 'none'}),
            ("As the price of a good declines, quantity demanded "
             "increases", 0),
            ("i.e., the demand curve is downward-sloping", 1),
        ],
        size=22, sub_size=20, line_spacing_pts=9)
    _add_hierarchical_bullets(
        slide, Inches(0.55), Inches(5.70), Inches(7.85), Inches(1.30),
        [
            ("Reasons: As the price declines...", 0,
             {'bold': True, 'bullet_style': 'none'}),
            ("More customers are willing to buy the product", 1),
            ("Existing customers buy more", 1),
        ],
        size=22, sub_size=20, line_spacing_pts=9)
    # native demand-curve mini figure (right column)
    fig = SimpleFig(9.35, 6.05, 3.0, 3.3, 10, 10)
    _fig_axes(slide, fig)
    # hand-tweaked from (fig.l-0.72, fig.b-fig.h-0.52) = (8.630, 2.230)
    # on 2026-08-23
    _add_text(slide, Inches(8.979), Inches(2.285),
              Inches(1.4), Inches(0.32), "Price", size=16, bold=True,
              italic=True, color=NAVY, font="Calibri")
    _add_text(slide, Inches(fig.l + fig.w - 0.55), Inches(fig.b + 0.10),
              Inches(1.5), Inches(0.32), "Quantity", size=16, bold=True,
              italic=True, color=NAVY, font="Calibri")
    _add_arrow(slide, (fig.x(0.8), fig.y(9.0)), (fig.x(9.2), fig.y(1.2)),
               color=NAVY, weight_pt=2.75, head=False)
    _add_text(slide, fig.x(9.35), fig.y(2.1), Inches(0.5), Inches(0.4),
              "D", size=20, bold=True, color=NAVY, font="Calibri")
    _draw_footer(slide, FOOTER_TEXT, page_num)
    return slide


# --------------------------------------------------------------------------
# Slide 10 — 1. MORE customers buy the product (Gjelina pizza, poll setup)
# --------------------------------------------------------------------------

def slide_10_more_customers(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_LAW)
    _draw_action_title(slide, "1. MORE Customers Buy the Product")
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.15), Inches(1.495), Inches(12.3),
        Inches(1.35),
        [
            ("How much are you willing to pay for a slice of (great) "
             "Gjelina pizza when you are hungry?", 0),
            ([("Maximum willingness to pay is your ", {}),
              ("Marginal Personal Value (MPV)", {'bold': True}),
              (" of the slice", {})], 0, {}),
        ],
        size=24, sub_size=22)
    # 2026-08-25 (Nico): his own cut-out, "Nice Pizza Slice.png",
    # replaces the wedge _mk_slice.py carved out of the Gjelina photo.
    # It has a transparent background and an irregular outline, so no
    # rounded corners - the drop shadow follows the alpha edge.
    _add_media_image(slide, "Nice Pizza Slice.png", left=Inches(4.765),
                     top=Inches(3.020), width=Inches(3.800),
                     rounded=False, shadow=True)
    _draw_footer(slide, FOOTER_TEXT, 10)
    _add_pollbreak_badge(slide)
    return slide


# --------------------------------------------------------------------------
# Slide 14 — 2. Existing customers buy MORE
# --------------------------------------------------------------------------

def slide_14_existing_buy_more(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_LAW)
    _draw_action_title(slide, "2. Existing Customers Buy MORE")
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.15), Inches(2.30), Inches(6.0),
        Inches(3.6),
        [
            ("What about a second slice of pizza?", 0),
            ("A third?", 1),
            ("A fourth?", 1),
            ("A millionth?", 1),
            ("… Diminishing Marginal Personal Value (MPV)!", 0),
        ],
        size=28, sub_size=24)
    _add_media_image(slide, "image14.jpeg", left=Inches(7.30),
                     top=Inches(1.95), width=Inches(4.45))
    _draw_footer(slide, FOOTER_TEXT, 14)
    return slide


# --------------------------------------------------------------------------
# Slide 15 — Law of Demand: multiple-unit consumers (native staircase;
# "MOV" typo fixed to "MPV" per 2026-08-14 decision)
# --------------------------------------------------------------------------

def slide_15_multiunit(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_LAW)
    _draw_action_title(slide,
                       "To the Same Consumer, Each Extra Slice is Worth Less")

    # figure lifted 0.20" by hand 2026-08-25 (was b = 6.05) to make
    # room for the boxed takeaway underneath
    fig = SimpleFig(2.6, 5.85, 7.6, 3.75, 5.8, 13.5)
    _fig_axes(slide, fig)
    _add_text(slide, Inches(2.00), Inches(1.52),
              Inches(1.6), Inches(0.32), "Price ($)", size=16, bold=True,
              italic=True, color=NAVY, font="Calibri")
    _add_text(slide, Inches(9.80), Inches(5.95),
              Inches(1.6), Inches(0.27), "Quantity (slices)", size=16,
              bold=True, italic=True, color=NAVY, font="Calibri")

    heights = [12, 9, 6, 3, 0]
    ordinals = ["1st", "2nd", "3rd", "4th", "5th"]
    bar_w_units = 0.55
    for i, (h, o) in enumerate(zip(heights, ordinals)):
        q = i + 1
        if h > 0:
            x_l = fig.x(q - bar_w_units / 2)
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x_l, fig.y(h),
                fig.x(q + bar_w_units / 2) - x_l, fig.y(0) - fig.y(h))
            bar.fill.solid()
            bar.fill.fore_color.rgb = GOLD
            bar.line.color.rgb = NAVY
            bar.line.width = Pt(1.25)
            bar.shadow.inherit = False
            _fig_ytick(slide, fig, h, "$%d" % h)
        _add_text(slide, fig.x(q + 0.30), fig.y(h) - Inches(0.30),
                  Inches(1.75), Inches(0.3), "MPV of %s slice" % o,
                  size=14, color=NAVY, font="Calibri")
        _fig_xtick(slide, fig, q, str(q))
    _fig_ytick(slide, fig, 0, "$0")
    # declining dashed guide over the bar tops
    _add_arrow(slide, (fig.x(1.0), fig.y(12.0)), (fig.x(5.0), fig.y(0.0)),
               color=NAVY, weight_pt=1.75, head=True, dash="dash")
    # Bottom takeaway, re-set by hand 2026-08-25: a narrow two-line
    # block, the second line 18 pt, sitting in a blue card.
    # repositioned by hand 2026-08-25 so the card clears the footer
    # rule (was 3.02 / 6.26, h 0.93)
    _add_rounded_filled_box(
        slide, Inches(2.85), Inches(6.20), Inches(6.45), Inches(0.80),
        "", fill=PALE_BLUE, text_color=NAVY, size=12,
        line=CBLUE, line_w=1.25, corner_pct=0.16, shadow=True)
    _add_hierarchical_bullets(
        slide, Inches(3.02), Inches(6.24), Inches(6.11), Inches(0.71),
        [([("Diminishing Marginal Personal Value  ", {'bold': True})], 0,
          {'bullet_style': 'none', 'align': PP_ALIGN.CENTER}),
         ([("(", {'size': 18}), ("MPV", {'size': 18}),
           (" declines as the ", {'size': 18}),
           ("same", {'size': 18, 'underline': True}),
           (" customer buys more)", {'size': 18})], 0,
          {'bullet_style': 'none', 'align': PP_ALIGN.CENTER,
           'space_before_pts': 0})],
        size=24)
    _set_notes(slide, (
        "Now that we have all the information needed, let’s draw your "
        "demand curve for pizza, step by step. For the first slice of pizza "
        "you are willing to pay 12 dollars. So if the price per slice was 12 "
        "dollars, you will be willing to buy one unit. The demand curve "
        "starts with this vertical bar at quantity 1 and a height of $12. "
        "You are willing to pay $9 for a second slice of pizza. So, if we "
        "reduce the price of the slice to $9, you will demand exactly two "
        "units of the good. We represent this information with this bar at "
        "Q equal to 2 and P equal to $9. For the third slice, you are "
        "willing to pay $6. If the restaurant offered the slices at $6 per "
        "piece, you’d buy exactly 3 slices. We represent this with a "
        "bar at Q equal to 3 and P equal to $6. For the fourth slice, you "
        "are willing to pay just $3. If the price per slice was $3, "
        "you’d buy exactly 4 slices. We represent this with a bar at Q "
        "equal to 4 and P equal to $3. You are willing to pay $0 for a "
        "fifth slice. If the restaurant charged $0 per slice, you’d buy "
        "5 slices. We represent this with a bar at Q equal to 5 and P equal "
        "to $0. Your demand function is like a staircase, listing how much "
        "you are willing to pay for the first slice of pizza, for the "
        "second slice, and so on and so forth. You can see that this demand "
        "function obeys the Law of Demand too. The reason for that is the "
        "Principle of Diminishing Marginal Utility. That is, the demand "
        "curve is downward sloping because you get satiated as you consume "
        "more and more pizza."))
    _draw_footer(slide, FOOTER_TEXT, 15)
    return slide


# --------------------------------------------------------------------------
# Slide 16 — Diminishing MPV (Bill Gates quote)
# --------------------------------------------------------------------------

def slide_16_gates(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_LAW)
    _draw_action_title(slide, "Diminishing Marginal Personal Value")
    # glyph repositioned by hand 2026-08-25 (was 0.85 / 2.10)
    _add_text(slide, Inches(1.38), Inches(1.83), Inches(1.2), Inches(1.9),
              "“", size=96, color=GRAY, font="Calibri")
    _add_hierarchical_bullets(
        slide, Inches(2.05), Inches(2.30), Inches(6.0), Inches(3.5),
        [
            ("I can understand wanting to have a million dollars… but "
             "once you get beyond that, I have to tell you, it’s the "
             "same hamburger.", 0,
             {'bullet_style': 'none', 'italic': True}),
            ("", 0, {'bullet_style': 'none'}),
            ("– Bill Gates", 0, {'bullet_style': 'none', 'bold': True,
                                 'size': 24}),
        ],
        size=28, line_spacing_pts=10)
    _add_media_image(slide, "image18.png", left=Inches(8.35),
                     top=Inches(2.35), width=Inches(4.15))
    _set_notes(slide, (
        "To further illustrate the principle of diminishing marginal "
        "utility, I’ll use a little help from Bill Gates. When asked "
        "what it feels like to be a billionaire, Bill Gates said “I can "
        "understand wanting to have a million dollars, but once you go "
        "beyond that, I have to tell you, it's the same hamburger.” In "
        "Economic terms, what Bill Gates is saying is that as your income "
        "increases your happiness goes up, but the marginal happiness that "
        "a dollar can buy is getting smaller and smaller. So once you make "
        "one million dollars per year, an additional dollar does not make "
        "any difference. Indeed, we can test Bill Gates’ hypothesis "
        "with some data. We use what social scientists call happiness data. "
        "This data is based on surveys, asking people how happy they are on "
        "a scale from 1 to 10, where 1 is very unhappy and 10 is very "
        "happy. I use happiness data in my own research. We can use the "
        "power of econometrics to learn what makes us happy. Do you want to "
        "know the secret to happiness? I would love to tell you, but you "
        "will have to pay for a Premium membership for that. I will show "
        "you one chart based on some happiness data. This graph shows the "
        "results from one specific study, but these findings are robust "
        "across hundreds of studies that have been conducted in almost "
        "every country."))
    _draw_footer(slide, FOOTER_TEXT, 16)
    return slide


# --------------------------------------------------------------------------
# Slide 17 — Income and life satisfaction (Inglehart 2018)
# --------------------------------------------------------------------------

def slide_17_inglehart(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_LAW)
    _draw_action_title(
        slide, "Diminishing MPV: Income and Life Satisfaction")
    pic_w = Inches(5.15)
    pic_x = int((SLIDE_W - pic_w) / 2)
    _add_media_image(slide, "image19.png", left=pic_x, top=Inches(1.75),
                     width=pic_w)
    _add_text(slide, pic_x, Inches(6.75), pic_w, Inches(0.3),
              "Source: Inglehart (2018)", size=12, italic=True, color=GRAY,
              font="Calibri", align=PP_ALIGN.CENTER)
    _set_notes(slide, (
        "The same diminishing pattern shows up far away from pizza. This "
        "is Inglehart's cross-country evidence on income and life "
        "satisfaction: satisfaction climbs steeply as income rises from "
        "low levels, then flattens out. Moving from very poor to "
        "comfortable changes a great deal; the next increment of income "
        "at an already high level changes much less. That is diminishing "
        "marginal personal value again, measured on a very different "
        "scale from a slice of pizza.\n\n"
        "Source: Cultural Evolution — People's Motivations are Changing, "
        "and Reshaping the World. Online publication date: March 2018, pp "
        "140-172. Ronald F. Inglehart, University of Michigan, Ann Arbor. "
        "https://doi.org/10.1017/9781108613880.009"))
    _draw_footer(slide, FOOTER_TEXT, 17)
    return slide


# --------------------------------------------------------------------------
# Slide 18 — Consumer Optimization (MB = MC recall)
# --------------------------------------------------------------------------

def slide_18_consumer_opt(prs):
    """Consumer optimization, redesigned 2026-08-25 on Nico's brief.

    The general rule and the consumption rule are the SAME box, drawn
    twice: navy for MB = MC (Module 1), gold for MPV = MC.  Between them
    sit the definitions, and a gold chevron carries the eye from one to
    the other.  Below the gold box, the two directions to move in.
    """
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_LAW)
    _draw_action_title(slide, "Consumer Optimization")

    # --- hero 1: the general rule, from Module 1 ---------------------
    _add_rounded_filled_box(
        slide, Inches(3.79), Inches(1.42), Inches(5.85), Inches(1.15),
        "", fill=NAVY, corner_pct=0.10, shadow=True)
    _add_hierarchical_bullets(
        slide, Inches(3.99), Inches(1.50), Inches(5.45), Inches(1.00),
        [("General optimization rule (from Module 1)", 0,
          {'bullet_style': 'none', 'align': PP_ALIGN.CENTER,
           'size': 17, 'bold': True, 'color': GOLD}),
         ("MB  =  MC", 0,
          {'bullet_style': 'none', 'align': PP_ALIGN.CENTER,
           'size': 32, 'bold': True, 'color': WHITE,
           'space_before_pts': 2})],
        size=17)

    # --- what MB and MC mean -----------------------------------------
    # tightened by hand 2026-08-25 (was 2.76 / 2.78, 8.65 x 1.70)
    _add_convention_box(
        slide, Inches(3.59), Inches(2.71), Inches(6.16), Inches(1.35),
        pad_h=Inches(0.14), pad_v=Inches(0.10),
        runs=[("MB: ", {'bold': True}),
              ("marginal benefit of an extra unit of a good/service",
               {}),
              ("MC: ", {'bold': True, 'newline': True}),
              ("marginal cost, which includes", {}),
              ("\u2013  the price of the good or service",
               {'newline': True, 'size': 16}),
              ("\u2013  the opportunity cost (e.g. time spent getting a "
               "haircut)", {'newline': True, 'size': 16})],
        size=18)

    # --- the eye travels from the general rule to the consumption one -
    _add_text(slide, Inches(6.17), Inches(4.12), Inches(1.0), Inches(0.30),
              "\u25bc", size=20, bold=True, color=GOLD, font="Calibri",
              align=PP_ALIGN.CENTER)

    # --- hero 2: the same rule, in the consumption context -----------
    _add_rounded_filled_box(
        slide, Inches(3.79), Inches(4.45), Inches(5.85), Inches(1.71),
        "", fill=GOLD, corner_pct=0.10, shadow=True)
    _add_hierarchical_bullets(
        slide, Inches(3.99), Inches(4.54), Inches(5.45), Inches(1.44),
        [("In the context of consumption we use MPV", 0,
          {'bullet_style': 'none', 'align': PP_ALIGN.CENTER,
           'size': 17, 'bold': True, 'color': NAVY}),
         ("(essentially the same as MB)", 0,
          {'bullet_style': 'none', 'align': PP_ALIGN.CENTER,
           'size': 13, 'color': NAVY, 'space_before_pts': 0}),
         ("Choose consumption quantity where:  MPV  =  MC", 0,
          {'bullet_style': 'none', 'align': PP_ALIGN.CENTER,
           'size': 24, 'bold': True, 'color': NAVY,
           'space_before_pts': 3})],
        size=17)

    # --- and which way to move ---------------------------------------
    _add_outlined_box(
        slide, Inches(3.79), Inches(6.38), Inches(2.80), Inches(0.56),
        "If MPV > MC :  buy more", line=NAVY, text_color=NAVY,
        size=16, bold=True, rounded=True, shadow=True, corner_pct=0.22)
    _add_outlined_box(
        slide, Inches(6.84), Inches(6.37), Inches(2.80), Inches(0.56),
        "If MPV < MC :  buy less", line=NAVY, text_color=NAVY,
        size=16, bold=True, rounded=True, shadow=True, corner_pct=0.22)

    # --- the recurring MB = MC anchor --------------------------------
    _add_anchor_burst(slide, Inches(11.41), Inches(5.28), Inches(1.69),
                      Inches(1.69), "MB = MC", top_size=20)
    _set_notes(slide, (
        "This is the same optimization rule as in Module 1, now applied "
        "to a consumer. One change in wording: instead of marginal "
        "benefit (MB) we write marginal personal value (MPV). It is the "
        "same concept – MPV is MB specific to consumption, the value the "
        "buyer places on one more unit. So “buy up to the point "
        "where MPV = MC” is the consumption version of MB = MC, "
        "which is why the anchor star sits next to the rule."))
    _draw_footer(slide, FOOTER_TEXT, 18)
    return slide


# --------------------------------------------------------------------------
# Slide 19 — Optimal number of movies (native MPV / MC chart)
# --------------------------------------------------------------------------

def slide_19_movies(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_LAW)
    _draw_action_title(
        slide, "Optimal Number of Movies to Watch "
               "(per Week, by One Customer)")
    fig = SimpleFig(2.1, 6.55, 9.0, 4.3, 10, 10)
    _fig_axes(slide, fig)
    _add_text(slide, Inches(fig.l - 1.3), Inches(fig.b - fig.h - 0.55),
              Inches(1.6), Inches(0.32), "$/movie", size=18, bold=True,
              italic=True, color=NAVY, font="Calibri")
    _add_text(slide, Inches(fig.l + fig.w - 1.1), Inches(fig.b + 0.10),
              Inches(2.4), Inches(0.32), "Movies per week", size=18,
              bold=True, italic=True, color=NAVY, font="Calibri")
    # MPV: declining convex curve (Bézier control points in fig units)
    _b0, _b1, _b2, _b3 = (0.35, 9.3), (2.2, 3.6), (4.6, 2.0), (9.4, 1.0)

    def _bez(t, a, b, c, d):
        return ((1 - t) ** 3 * a + 3 * (1 - t) ** 2 * t * b
                + 3 * (1 - t) * t * t * c + t ** 3 * d)

    _mpv_curve = _add_cubic_curve(
        slide, (fig.x(_b0[0]), fig.y(_b0[1])),
        (fig.x(_b1[0]), fig.y(_b1[1])),
        (fig.x(_b2[0]), fig.y(_b2[1])),
        (fig.x(_b3[0]), fig.y(_b3[1])),
        color=NAVY, weight_pt=4.0)
    # soft shade on the MPV curve (Nico, 2026-08-23)
    _add_drop_shadow(_mpv_curve, blur="38100", dist="25400",
                     direction="2700000", alpha="40000")
    # label hand-tweaked from (fig.x(8.6), fig.y(0.5)-0.42) = (9.840,
    # 5.915) on 2026-08-23
    _add_text(slide, Inches(10.700), Inches(6.059), Inches(1.0),
              Inches(0.4), "MPV", size=20, bold=True, italic=True,
              color=NAVY, font="Calibri")
    # MC line (red).  UPWARD-SLOPING because the label says the curve
    # includes opportunity cost: each extra movie crowds out the next-best
    # use of an hour, and the best alternatives go first (Nico's hand
    # redraw, 2026-08-23 — now a standing rule in the Teaching CLAUDE.md).
    _mc0, _mc1 = (0.4173, 0.4884), (8.3778, 4.3360)
    _add_arrow(slide, (fig.x(_mc0[0]), fig.y(_mc0[1])),
               (fig.x(_mc1[0]), fig.y(_mc1[1])),
               color=RED, weight_pt=2.5, head=False)
    # label hand-placed above the right end of the MC line (2026-08-23)
    _add_text(slide, Inches(8.676), Inches(4.322), Inches(3.6),
              Inches(0.4), "MC (incl. opportunity cost)", size=18,
              italic=True, color=RED, font="Calibri")
    # Q* drop line at the TRUE MPV/MC crossing — solve the Bézier against
    # the sloped MC line (the marked Q* must sit at the actual
    # intersection, per the "curves must be economically exact" rule)
    _mc_slope = (_mc1[1] - _mc0[1]) / (_mc1[0] - _mc0[0])

    def _gap(t):
        xs = _bez(t, _b0[0], _b1[0], _b2[0], _b3[0])
        ys = _bez(t, _b0[1], _b1[1], _b2[1], _b3[1])
        return ys - (_mc0[1] + _mc_slope * (xs - _mc0[0]))

    lo, hi = 0.0, 1.0
    for _ in range(80):
        mid = (lo + hi) / 2
        if _gap(mid) > 0:
            lo = mid
        else:
            hi = mid
    t_star = (lo + hi) / 2
    qstar = _bez(t_star, _b0[0], _b1[0], _b2[0], _b3[0])
    vstar = _bez(t_star, _b0[1], _b1[1], _b2[1], _b3[1])
    _add_arrow(slide, (fig.x(qstar), fig.y(vstar)), (fig.x(qstar), fig.y(0)),
               color=NAVY, weight_pt=1.5, head=False, dash="dash")
    _add_text(slide, fig.x(qstar) - Inches(0.3), Inches(fig.b + 0.08),
              Inches(0.6), Inches(0.35), "Q*", size=20, bold=True,
              italic=True, color=NAVY, font="Calibri",
              align=PP_ALIGN.CENTER)
    # annotations
    _add_text(slide, Inches(5.7), Inches(1.85), Inches(6.9), Inches(0.4),
              "MPV diminishes as more and more movies have been watched",
              size=18, color=NAVY, font="Calibri")
    # box resized + moved up by hand on 2026-08-23 (was 5.55, 3.05,
    # 4.4 x 1.15 with the default 0.20/0.12 text padding)
    _add_convention_box(
        slide, Inches(5.500), Inches(2.870), Inches(3.800), Inches(1.028),
        pad_h=Inches(0.173), pad_v=Inches(0.108),
        # "is" bolded as well as underlined (Nico, 2026-08-23)
        runs=[("MPV ", {'bold': True}),
              ("is", {'bold': True, 'underline': True}),
              (" the demand curve ", {'bold': True}),
              ("(strictly speaking, the “inverse” demand curve, "
               "but economists are not very precise about this)",
               {'size': 13})],
        size=16)
    # arrow ends ON the MPV curve (point evaluated on the Bézier)
    t_a = 0.30
    _ax = _bez(t_a, _b0[0], _b1[0], _b2[0], _b3[0])
    _ay = _bez(t_a, _b0[1], _b1[1], _b2[1], _b3[1])
    # 2026-08-25 (Nico): the pointer to the MPV curve takes the curve's
    # own dark blue and a heavier weight
    _add_arrow(slide, (Inches(5.50), Inches(3.85)),
               (fig.x(_ax), fig.y(_ay)), color=NAVY, weight_pt=3.0,
               head=True)
    _draw_footer(slide, FOOTER_TEXT, 19)
    return slide


# --------------------------------------------------------------------------
# Slide 20 — Aggregating individual consumers' demand
# (clean-number rebuild: C1 and C2 dots sum horizontally to aggregate)
# --------------------------------------------------------------------------

def slide_20_aggregation(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_LAW)
    _draw_action_title(slide, "Aggregating Individual Consumers’ Demand")
    fig = SimpleFig(1.85, 6.45, 10.4, 4.15, 12.5, 14.5)
    _fig_axes(slide, fig)
    _add_text(slide, Inches(fig.l - 1.35), Inches(fig.b - fig.h - 0.55),
              Inches(1.7), Inches(0.32), "Price ($)", size=18, bold=True,
              italic=True, color=NAVY, font="Calibri")
    _add_text(slide, Inches(fig.l + fig.w - 0.7), Inches(fig.b + 0.10),
              Inches(1.7), Inches(0.32), "Quantity", size=18, bold=True,
              italic=True, color=NAVY, font="Calibri")
    for p in (12, 9, 6, 3):
        _fig_ytick(slide, fig, p, str(p), size=18)
    for q in range(1, 12):
        _fig_xtick(slide, fig, q, str(q), size=18)
    c1 = [(1, 12), (2, 9), (3, 6), (4, 3)]
    c2 = [(2, 12), (3, 9), (4, 6), (5, 3)]
    agg = [(3, 12), (5, 9), (7, 6), (9, 3)]
    # demand lines through the dots (extended slightly beyond)
    _add_arrow(slide, (fig.x(0.65), fig.y(13.05)), (fig.x(4.7), fig.y(0.9)),
               color=NAVY, weight_pt=2.0, head=False)
    _add_arrow(slide, (fig.x(1.65), fig.y(13.05)), (fig.x(5.7), fig.y(0.9)),
               color=GOLD, weight_pt=2.0, head=False)
    _add_arrow(slide, (fig.x(2.3), fig.y(13.05)), (fig.x(9.7), fig.y(1.9)),
               color=GREEN, weight_pt=2.5, head=False)
    for (q1, p), (q2, _p2), (qa, _pa) in zip(c1, c2, agg):
        _fig_point(slide, fig, q1, p, fill=NAVY)
        _fig_point(slide, fig, q2, p, fill=GOLD)
        _fig_point(slide, fig, qa, p, fill=GREEN)
        _add_text(slide, fig.x((q1 + q2) / 2) - Inches(0.14),
                  fig.y(p) - Inches(0.16), Inches(0.3), Inches(0.3),
                  "+", size=16, bold=True, color=GRAY, font="Calibri",
                  align=PP_ALIGN.CENTER)
        _add_text(slide, fig.x((q2 + qa) / 2) - Inches(0.14),
                  fig.y(p) - Inches(0.16), Inches(0.3), Inches(0.3),
                  "=", size=16, bold=True, color=GRAY, font="Calibri",
                  align=PP_ALIGN.CENTER)
    # legend (top-right, stacked)
    ly = Inches(1.95)
    for label, color in (("Consumer 1 demand", NAVY),
                         ("Consumer 2 demand", GOLD),
                         ("Aggregate demand", GREEN)):
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.55),
                                     int(ly + Inches(0.08)), Inches(0.16),
                                     Inches(0.16))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()
        dot.shadow.inherit = False
        _add_text(slide, Inches(9.82), ly, Inches(3.2), Inches(0.35),
                  label, size=18, bold=(color is GREEN), color=color,
                  font="Calibri")
        ly = int(ly + Inches(0.42))
    _draw_footer(slide, FOOTER_TEXT, 20)
    return slide




# ==========================================================================
#  BATCH C — Part 1 close (21–24) + elasticities core (26–45)
# ==========================================================================

# OMML shortcuts used throughout the elasticity section
def _oED():
    return _omml_sub(_omml_run('E'), _omml_text('D'))


def _oEI():
    return _omml_sub(_omml_run('E'), _omml_run('I'))


def _oEXY():
    return _omml_sub(_omml_run('E'),
                     _omml_run('X') + _omml_text(',') + _omml_run('Y'))


def _o_pct(var):
    return _omml_text('%Δ') + _omml_run(var)


def _oED_frac(x_var='P'):
    return (_oED() + _omml_text(' = ')
            + _omml_frac(_o_pct('Q'), _o_pct(x_var)))


def _custom_title_runs(slide, runs, *, size=30):
    """Action title with per-run colors (e.g. red 'Method 1')."""
    _add_rect(slide, MARGIN, Inches(1.28), RULE_W, Inches(0.02), RULE)
    _add_rect(slide, MARGIN, Inches(1.265), GOLD_W, Inches(0.05), GOLD)
    _add_hierarchical_bullets(
        slide, MARGIN, Inches(0.55), RULE_W, Inches(0.7),
        [(runs, 0, {'bullet_style': 'none', 'bold': True, 'size': size,
                    'color': NAVY})],
        size=size)


# --------------------------------------------------------------------------
# Slide 21 — Factors affecting demand (+ CT demand-shift mini graphic)
# --------------------------------------------------------------------------

def slide_21_factors(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_LAW)
    _draw_action_title(slide, "Factors Affecting Demand")
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.15), Inches(1.90), Inches(7.6),
        Inches(5.0),
        [
            ("Income", 0),
            ("Prices of related goods", 0),
            ("Complements versus substitutes", 1),
            ("Complements can be used strategically", 2,
             {'bullet_style': 'sub', 'size': 20}),
            ("E.g., Tesla & Superchargers. Ryanair.", 2,
             {'bullet_style': 'sub', 'size': 20}),
            ("Price expectations", 0),
            ("Population size and composition", 0),
            ([("Product quality, advertising, placement ", {}),
              ("(within the firm’s control)",
               {'italic': True, 'size': 18})], 0, {}),
            ("Network effects", 0),
        ],
        size=24, sub_size=22, line_spacing_pts=12)
    # cartoon moved to the lower-right corner (Nico, 2026-08-23; was
    # 9.15, 1.65)
    _add_media_image(slide, "image20.png", left=Inches(10.400),
                     top=Inches(3.831), width=Inches(2.35))
    # mini demand-shift figure (CT adoption): D shifts right / left.
    # Nico redrew it by hand on 2026-08-23 — the figure moved up beside
    # the bullets (was anchored at 9.15, 6.75), the two shifted curves
    # each got a short diagonal arrow off the base curve, and the labels
    # became "Rising demand" / "Falling demand" beside their own curve.
    # Each shifted curve + its arrow + its label is ONE group (see
    # MANUAL_GROUPS in _group_pass.py), revealed on one click.
    fig = SimpleFig(6.876, 4.003, 2.6, 2.2, 10, 10)
    _fig_axes(slide, fig, weight_pt=1.5)
    _add_arrow(slide, (fig.x(1.4769), fig.y(9.0)), (fig.x(8.2), fig.y(1.6)),
               color=NAVY, weight_pt=2.25, head=False)
    # right-shift panel: dashed curve, arrow, label (emitted together so
    # the grouping pass finds them adjacent and in this order)
    _add_arrow(slide, (fig.x(3.5192), fig.y(9.4682)),
               (fig.x(9.7538), fig.y(3.2636)),
               color=GRAY, weight_pt=1.5, head=False, dash="dash")
    _add_arrow(slide, (fig.x(5.3), fig.y(5.0136)),
               (fig.x(6.4654), fig.y(6.1727)),
               color=GOLD, weight_pt=2.0, head=True)
    _add_text(slide, Inches(9.226), Inches(3.395), Inches(1.600),
              Inches(0.185), "Rising demand", size=11,
              italic=True, color=GRAY, font="Calibri")
    # left-shift panel
    _add_arrow(slide, (fig.x(0.2), fig.y(7.6)), (fig.x(6.2), fig.y(0.6)),
               color=GRAY, weight_pt=1.5, head=False, dash="dash")
    _add_arrow(slide, (fig.x(4.3231), fig.y(5.1)),
               (fig.x(3.2), fig.y(4.1)),
               color=GOLD, weight_pt=2.0, head=True)
    _add_text(slide, Inches(7.188), Inches(3.653), Inches(1.126),
              Inches(0.185), "Falling demand", size=11,
              italic=True, color=GRAY, font="Calibri",
              align=PP_ALIGN.RIGHT)
    # moved right by hand on 2026-08-23 (was 3.05, 6.42)
    _add_outlined_box(slide, Inches(4.066), Inches(6.473), Inches(2.6),
                      Inches(0.48), "Anything else?", line=NAVY,
                      text_color=NAVY, size=16, bold=True, rounded=True,
                      shadow=True, corner_pct=0.25)
    _set_notes(slide, (
        "INCOME (or wealth): When income goes up, demand for normal goods "
        "goes up, e.g. vacations. Inferior goods are such that demand goes "
        "down when the income goes up, e.g.: Ramen Noodles. Store brands "
        "tend to be inferior goods. It is helpful for a company to know "
        "whether they are selling a normal or an inferior good. Even "
        "better, knowing the precise quantitative relationship between "
        "personal income or wealth and demand would be nice. Example: "
        "Single malt whisky. As income in China goes up, prices of single "
        "malts go up by 10-15% a year. And that will continue because "
        "supply is constrained over a long period of time (supply on the "
        "market now may have been put in barrels as long ago as 30 years). "
        "PRICE OF RELATED GOODS: The main distinction here is between "
        "substitutes and complementary goods. Complements: Windows and "
        "Office. iTunes and iPhone. Goods such that consuming one makes "
        "your utility from consuming the other rise. Another example is "
        "telephone service in hotels. This used to be a significant source "
        "of revenue for the hospitality industry before the advent of cell "
        "phones. In 5 star hotels at one point, ancillary goods contributed "
        "up to 20% of the revenues from one room (room service, phone "
        "service, laundry, etc.) on average. So there were incentives to "
        "drive down the price of the room to drive the consumption of "
        "ancillary services. Now that cell phones have emerged the "
        "incentive to drive down the price of hotel rooms to support "
        "demand for complements/ancillaries has diminished. The price of "
        "the ancillaries themselves has risen steeply. In the computer "
        "industry, declining prices of hardware has supported demand for "
        "complements like software. Microsoft has benefitted a lot from "
        "the long term decline in the price of PCs. Here, an upward price "
        "change along a complementary good's demand curve shifts demand "
        "for the other good inward. Substitutes: Frozen yogurt versus ice "
        "cream. Butter versus margarine. Lexus versus Mercedes. Used cars "
        "versus new cars. Cable internet access versus DSL. Goods such "
        "that consuming one makes it less likely I want to consume the "
        "other. As the price of DSL falls, the demand for cable access "
        "falls. As the price of butter rises, the demand for margarine "
        "rises. Here, an upward price change along a good's demand curve "
        "shifts demand for the other good outward. PRICE EXPECTATIONS: If "
        "you believe the price will rise, you want to stock up. This can "
        "lead to a self-fulfilling prophecy as my increase in demand "
        "drives up the price. This can be an important phenomenon on "
        "financial markets, and can explain bubbles based on beliefs about "
        "the evolution of future prices. POPULATION SIZE AND COMPOSITION: "
        "Probably the most underrated of the determinants. When population "
        "rises fast, you can anticipate selling more. But the structure of "
        "the population might be even more important. A younger population "
        "might be more attracted to music and movies. This can and must "
        "play into your business plan. A big advantage of demographics is "
        "predictability. A 25 year old today will be 26 next year with "
        "near certainty. Example: Nestlé was looking to buy Ortega Foods a "
        "few years ago, and did. They were anticipating a rise in the "
        "demand from Hispanics for the types of products Ortega sold. They "
        "thought demand would rise, which made Ortega an attractive "
        "acquisition target. NETWORK EFFECTS: These used to be called snob "
        "and bandwagon effects. Snob effect: I might have a snobbish "
        "attitude toward seeing so many people driving a similar car to "
        "mine, and buy something else. The existence of high levels of "
        "demand for a product creates a disincentive effect in purchase "
        "decisions. The bandwagon effect is the opposite. A friend has a "
        "nice blue shirt, which makes me want it too. Demand by others "
        "affects my propensity to demand the good, either negatively "
        "(snob) or positively (bandwagon). In modern economics these "
        "effects are called network effects."))
    _draw_footer(slide, FOOTER_TEXT, 21)
    return slide


# --------------------------------------------------------------------------
# Slide 22 — The Snob Effect in the News (NEW, from CT; WSJ clippings)
# --------------------------------------------------------------------------

def slide_22_snob_news(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_LAW)
    _draw_action_title(slide, "The Snob Effect in the News")
    # one-line reminder of the concept above the two clippings (Nico,
    # 2026-08-23); both panels shifted down 0.20" to make room
    # width measured with PIL on Calibri/Calibri Bold at 18 pt (9.11" of
    # text, so 9.50" of text box) — the line must not wrap, and the box
    # has to stay under the 10" "layout band" threshold that makes
    # _group_pass rule 1 skip a filled roundRect
    _add_convention_box(
        slide, Inches(1.717), Inches(1.400), Inches(9.900), Inches(0.520),
        runs=[("Snob effect: ", {'bold': True}),
              ("exclusivity is part of the value, so demand falls as "
               "more people own the good", {})],
        size=18, align=PP_ALIGN.CENTER)
    # left panel: Ferrari (headline above photo)
    # panels nudged down by hand 2026-08-25 (were 2.05 / 3.35 and
    # 2.15 / 3.75)
    _add_media_image(slide, "ct_snob_image7.png", left=Inches(0.75),
                     top=Inches(2.51), width=Inches(5.6),
                     rounded=False, shadow=False)
    _add_media_image(slide, "ct_snob_image6.jpg", left=Inches(1.30),
                     top=Inches(3.81), width=Inches(4.5))
    # right panel: Birkin (headline above photo)
    _add_media_image(slide, "ct_snob_image9.png", left=Inches(7.10),
                     top=Inches(2.21), width=Inches(5.4),
                     rounded=False, shadow=False)
    _add_media_image(slide, "ct_snob_image8.jpg", left=Inches(7.75),
                     top=Inches(3.81), width=Inches(4.1))
    _add_text(slide, MARGIN, Inches(6.78), Inches(6.0), Inches(0.3),
              "Source: The Wall Street Journal (2025)", size=12,
              italic=True, color=GRAY, font="Calibri")
    _set_notes(slide, (
        "Two recent examples of the snob effect at work. Ferrari caps "
        "production well below demand — you cannot simply buy its top "
        "models, you have to be invited — and that engineered scarcity "
        "has made it the most valuable car company in Europe. Hermès "
        "does the same with the Birkin bag. Exclusivity itself is what "
        "customers are paying for: high demand by others would reduce "
        "the appeal."))
    _draw_footer(slide, FOOTER_TEXT, 22)
    return slide


# --------------------------------------------------------------------------
# Slide 23 — Implications of positive network effects
# --------------------------------------------------------------------------

def slide_23_network_effects(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_LAW)
    _draw_action_title(slide, "Implications of Positive Network Effects")
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.15), Inches(1.70), Inches(12.3),
        Inches(1.55),
        [
            ("Switching more difficult for customers", 0),
            ("“Tippy” Markets", 0),
            ("Winner-take-all", 0),
        ],
        size=26, line_spacing_pts=10)
    # nudged right and up by hand on 2026-08-23 (was 1.65, 3.35)
    _add_media_image(slide, "image21.jpg", left=Inches(3.200),
                     top=Inches(3.050), width=Inches(10.0))
    _draw_footer(slide, FOOTER_TEXT, 23)
    return slide


# --------------------------------------------------------------------------
# Slide 24 — Remember about demand
# --------------------------------------------------------------------------

def slide_24_remember(prs):
    # 2026-08-25 (Nico): the bullets move to a narrower left column and
    # his "Demand Shopping" photo fills the right side, so the slide
    # carries a picture rather than a full-width wall of text
    slide = make_content_bulleted(
        prs, 24, TAG_LAW, "Remember About Demand",
        [
            ("Represents the willingness to pay by all actual or "
             "potential customers for a good/service", 0),
            ("Price changes while everything else is kept constant", 0),
            ("We usually represent demand as a line, for simplicity", 0),
            ([("We draw ", {}), ("P", {'italic': True}),
              (" on the vertical axis, but mathematically:", {})], 0, {}),
            ([("Q", {'italic': True}), (" = 10 − ", {}),
              ("P", {'italic': True}),
              (" : this is the demand function (", {}),
              ("Q", {'italic': True}), (" as function of ", {}),
              ("P", {'italic': True}), (")", {})], 1, {}),
            ([("P", {'italic': True}), (" = 10 − ", {}),
              ("Q", {'italic': True}),
              (" : is the “inverse” demand function ", {}),
              ("→ so strictly speaking, we draw the “inverse” "
               "demand function", {'size': 20})], 1, {}),
        ],
        # column narrowed and lifted by hand 2026-08-25
        size=24, sub_size=22, bullets_width=Inches(7.31),
        bullets_top=Inches(1.62))
    # enlarged and lifted by hand 2026-08-25 (was 8.90 / 2.35, w 3.95)
    _add_media_image(slide, "Demand Shopping.png", left=Inches(8.45),
                     top=Inches(1.66), width=Inches(4.61))
    return slide


# --------------------------------------------------------------------------
# Slide 26 — Generic definition of elasticity
# --------------------------------------------------------------------------

def slide_26_generic_elasticity(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_ELAST)
    _draw_action_title(slide, "Generic Definition of Elasticity")
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.15), Inches(2.35), Inches(7.7),
        Inches(2.6),
        [
            ([("Refers to the ", {}),
              ("responsiveness to external forces", {'bold': True}),
              (":", {})], 0, {}),
            ([("Elastic", {'bold': True, 'underline': True}),
              (": changes when punched", {})], 1, {}),
            ([("Inelastic", {'bold': True, 'underline': True}),
              (": stays in original position when punched", {})], 1, {}),
        ],
        size=26, sub_size=24)
    _add_media_image(slide, "image22.png", left=Inches(8.95),
                     top=Inches(2.30), width=Inches(3.0))
    _set_notes(slide, (
        "In Economics, elasticities come in all shapes and sizes. There "
        "are demand elasticities and supply elasticities. The own-price "
        "elasticities and the cross-price elasticities. Elasticity is "
        "kind of a weird term, right? Economists borrowed the word "
        "elasticity from physics. In the early 20th century, the "
        "Economics discipline was going through a revolution. Many "
        "economists felt like physics was a real science, but economics "
        "wasn't. So in order to make economics a real science, they felt "
        "they needed to talk like physicists. That is why the Economic "
        "profession borrowed so many terms from Physics, such as "
        "elasticity and equilibrium. In physics, elasticity refers to the "
        "responsiveness to external forces. If an object is elastic, that "
        "means that if you punch it, it will react by adopting a "
        "different shape. If an object is not elastic, when you punch it, "
        "it will conserve its original form, it will not adapt. So, you "
        "can think of elasticity like the ability to adapt to external "
        "forces."))
    _draw_footer(slide, FOOTER_TEXT, 26)
    return slide


# --------------------------------------------------------------------------
# Slide 27 — Netflix price increases (native price-history chart, CT item 4)
# --------------------------------------------------------------------------

# US Standard plan, list price at each increase.  The March 2026 step
# to $19.99 verified 2026-08-25 against CNBC / Variety / CBS.
NETFLIX_PRICES = [(2010, 7.99), (2014, 8.99), (2015, 9.99), (2017, 10.99),
                  (2019, 12.99), (2020, 13.99), (2022, 15.49),
                  (2025, 17.99), (2026, 19.99)]


def slide_27_netflix(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_ELAST)
    _draw_action_title(
        slide, "We Use Elasticity to Study Effects of Price Increases")
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.15), Inches(2.00), Inches(6.3),
        Inches(4.6),
        [
            ("Netflix frequently raises its monthly subscription price",
             0, {'bold': True}),
            ("What are the implications of price increases for:", 0,
             {'bold': True}),
            ("Number of subscriptions?", 1),
            ("Revenues?", 1),
            ("Key concept: Elasticity", 0, {'bold': True, 'size': 28}),
        ],
        size=24, sub_size=22)
    # 2026-08-25 (Nico): the practice-video pointer is a reference box
    # now, not a bullet
    _add_reference_box(
        slide, Inches(0.60), Inches(5.62), Inches(5.85), Inches(0.78),
        "Module 2, Practice Video 1", kind="video", size=16,
        sub_label="numerical example with Netflix price increase",
        sub_size=13)
    # native chart: U.S. Standard-plan monthly price
    _add_media_image(slide, "image24.png", left=Inches(8.55),
                     top=Inches(1.70), width=Inches(1.7),
                     rounded=False, shadow=False)
    steps = []
    for i, (yr, pr) in enumerate(NETFLIX_PRICES):
        if i:
            steps.append((yr, NETFLIX_PRICES[i - 1][1]))
        steps.append((yr, pr))
    # run the line a little past the last step and give the axis
    # headroom, so the $19.99 label has room in the top-right corner
    steps.append((2027, NETFLIX_PRICES[-1][1]))
    _make_xy_line_chart(
        slide, Inches(7.15), Inches(2.60), Inches(5.55), Inches(3.55),
        series=[("Standard plan", steps, NAVY, 'none')],
        x_title="Year", y_title="$ / month",
        x_min=2009, x_max=2028, x_unit=4, y_min=0, y_max=22, y_unit=5)
    # gold callout on the 2014 hike — the Practice Video 1 example
    # callout, arrow and the two price labels repositioned by hand
    # 2026-08-25
    _add_text(slide, Inches(9.18), Inches(4.63), Inches(2.3), Inches(0.6),
              "+$1 in 2014\n($7.99 → $8.99)", size=12, bold=True,
              italic=True, color=GOLD, font="Calibri")
    # arrow and price label repositioned by hand 2026-08-25
    _add_arrow(slide, (Inches(9.279), Inches(4.612)),
               (Inches(9.070), Inches(4.383)), color=GOLD, weight_pt=1.75,
               head=True)
    _add_text(slide, Inches(8.05), Inches(4.52), Inches(1.4), Inches(0.3),
              "$7.99", size=14, bold=True, color=NAVY, font="Calibri")
    # 2026-08-25 (Nico): the current price, in the top-right corner
    _add_text(slide, Inches(11.05), Inches(2.62), Inches(1.6), Inches(0.3),
              "$19.99", size=15, bold=True, color=NAVY, font="Calibri",
              align=PP_ALIGN.RIGHT)
    _add_text(slide, Inches(7.15), Inches(6.25), Inches(5.55), Inches(0.3),
              "Netflix U.S. Standard plan · Source: Netflix price "
              "announcements", size=11, italic=True, color=GRAY,
              font="Calibri", align=PP_ALIGN.CENTER)
    _set_notes(slide, (
        "Why we need elasticity at all. Netflix raises its U.S. standard "
        "plan price regularly – the chart tracks it from $7.99 to "
        "$17.99. Every one of those increases faced the same two "
        "questions: how many subscribers will we lose, and will revenue "
        "go up or down? Notice those are different questions with "
        "different answers, and elasticity is the single number that "
        "links them. The revenue side is worked through in Module 2, "
        "Practice Video 1.\n\n"
        "Will be answered in Practice Video 1 for Module 2"))
    _draw_footer(slide, FOOTER_TEXT, 27)
    return slide


# --------------------------------------------------------------------------
# Slide 28 — What is the elasticity of demand?
# --------------------------------------------------------------------------

def slide_28_what_is_elasticity(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_ELAST)
    _draw_action_title(slide, "What Is the Elasticity of Demand?")
    _add_math_equation(
        slide, Inches(4.55), Inches(1.70), Inches(4.2), Inches(1.15),
        _omml_run('E') + _omml_text(' = ')
        + _omml_frac(_o_pct('Q'), _o_pct('X')),
        size_pt=30, color=NAVY, fill=CREAM, line=NAVY, rounded=True,
        shadow=True)
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.15), Inches(3.30), Inches(12.4),
        Inches(3.3),
        [
            ([("Definition:", {'bold': True}),
              (" How quantity demanded (", {}), ("Q", {'italic': True}),
              (") changes with respect to another variable (", {}),
              ("X", {'italic': True}), (")", {})], 0, {}),
            ([("E.g., ", {}), ("X", {'italic': True}),
              (" = price, income…", {})], 1, {}),
            ([("Characteristic:", {'bold': True}), ("  ", {}),
              ("Unit free", {'color': RED})], 0, {}),
            ([("How to talk about elasticity:", {'bold': True}),
              (" A % change in quantity demanded (", {}),
              ("Q", {'italic': True}),
              (") due to a 1% change in another variable (", {}),
              ("X", {'italic': True}), (")", {})], 0, {}),
        ],
        size=24, sub_size=22)
    _set_notes(slide, (
        "Elasticity can be a difficult concept to grasp, but we will work "
        "with it so much that you will become very familiar with the "
        "concept. Elasticity is the responsiveness or sensitivity of some "
        "variable to the change in some other variable. It is a unitless "
        "ratio of two percentage changes. Elasticities have no units, "
        "which make them directly comparable to each other. We express "
        "the changes in the variables as a percentage so that the "
        "calculated elasticity is just a number. Example: You can "
        "directly compare the elasticity of demand for Sara Lee products "
        "to the elasticity of demand for gasoline."))
    _draw_footer(slide, FOOTER_TEXT, 28)
    # 2026-08-25 (Nico, copied from CT): the numerator's percentage
    # symbol is circled, and the label points at the circle.  The three
    # pieces are grouped in _group_pass so they move as one.
    _add_text(slide, Inches(8.10), Inches(1.32), Inches(2.6), Inches(0.35),
              "percentage change", size=15, italic=True, color=GRAY,
              font="Calibri", align=PP_ALIGN.LEFT)
    _add_arrow(slide, (Inches(8.050), Inches(1.430)),
               (Inches(7.200), Inches(1.790)), color=GOLD, weight_pt=1.75,
               head=True)
    _add_oval_outline(slide, Inches(6.625), Inches(1.745), Inches(0.664),
                      Inches(0.535), color=GOLD, weight_pt=2.25)
    return slide


# --------------------------------------------------------------------------
# Slide 29 — Three types of elasticity (NEW, from CT)
# --------------------------------------------------------------------------

def _three_types_cards(slide, *, dim=None):
    """Three formula cards side by side; dim = set of indices to fade."""
    cards = [
        ("Own-price elasticity", _oED_frac('P')),
        ("Income elasticity",
         _oEI() + _omml_text(' = ')
         + _omml_frac(_o_pct('Q'), _omml_text('%Δ') + _omml_run('I'))),
        ("Cross-price elasticity",
         _oEXY() + _omml_text(' = ')
         + _omml_frac(_omml_text('%Δ')
                      + _omml_sub(_omml_run('Q'), _omml_run('X')),
                      _omml_text('%Δ')
                      + _omml_sub(_omml_run('P'), _omml_run('Y')))),
    ]
    xs = [Inches(0.55), Inches(4.85), Inches(9.15)]
    for i, (label, omml) in enumerate(cards):
        faded = dim is not None and i in dim
        color = FADED if faded else NAVY
        _add_rounded_filled_box(
            slide, xs[i], Inches(2.35), Inches(3.6), Inches(0.62), label,
            fill=FADED if faded else NAVY, text_color=WHITE, size=19,
            bold=True, corner_pct=0.12, shadow=not faded)
        _add_math_equation(
            slide, xs[i], Inches(3.30), Inches(3.6), Inches(1.5), omml,
            size_pt=26, color=color, fill=CREAM if not faded else None,
            line=color, rounded=True, shadow=not faded)


def slide_29_three_types(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_ELAST)
    _draw_action_title(slide, "Three Types of Elasticity")
    _three_types_cards(slide)
    _add_text(slide, MARGIN, Inches(5.45), RULE_W, Inches(0.45),
              "All three follow the same recipe: a % change in quantity, "
              "divided by the % change in the driving variable",
              size=20, italic=True, color=GRAY, font="Calibri",
              align=PP_ALIGN.CENTER)
    _set_notes(slide, (
        "Before we work through each one, here is the map. All "
        "elasticities in this module have the same structure: the "
        "percentage change in quantity demanded, divided by the "
        "percentage change in whatever is driving it — the good's own "
        "price, the customer's income, or the price of another good. We "
        "start with the own-price elasticity."))
    _draw_footer(slide, FOOTER_TEXT, 29)
    return slide


# --------------------------------------------------------------------------
# Slide 30 — Own-price elasticity (definition slide)
# --------------------------------------------------------------------------

def slide_30_own_price(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_OWN)
    _draw_action_title(slide, "Own-Price Elasticity")
    _add_math_equation(
        slide, Inches(4.35), Inches(1.75), Inches(4.6), Inches(1.35),
        _oED_frac('P'), size_pt=34, color=NAVY, fill=CREAM, line=NAVY,
        rounded=True, shadow=True)
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.15), Inches(3.60), Inches(12.4),
        Inches(3.0),
        [
            ("% change in quantity demanded, divided by the % change in "
             "price (of the same product)", 0),
            ([("Intuition:", {'bold': True}),
              (" How sensitive demand is to changes in the own price of "
               "the good", {})], 0, {}),
            ([("Sign of ", {}), ("E", {'italic': True}),
              ("D", {'italic': True, 'size': 16}),
              (" :  ", {}), ("Negative", {'bold': True}),
              ("   (Law of Demand:  ", {}), ("P", {'italic': True}),
              (" ↑  ⇒  ", {}), ("Q", {'italic': True}),
              (" ↓ )", {})], 0, {}),
        ],
        size=26, sub_size=24)
    _draw_footer(slide, FOOTER_TEXT, 30)
    return slide


# --------------------------------------------------------------------------
# Slide 31 — Example: demand elasticity for water (SmartArt → native boxes)
# --------------------------------------------------------------------------

def slide_31_water(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_OWN)
    _draw_action_title(slide, "Example: Demand Elasticity for Water")
    _add_media_image(slide, "image27.png", left=Inches(0.70),
                     top=Inches(2.00), width=Inches(4.0))
    steps = [
        ("Consulting company estimated: price elasticity of demand for "
         "water in L.A. is −0.4", NAVY, WHITE),
        ("LADWP wants to reduce quantity demanded by 10%", NAVY, WHITE),
        ("By how much should it raise its price?", GOLD, NAVY),
    ]
    y = Inches(2.05)
    for text, fill, tcol in steps:
        _add_rounded_filled_box(slide, Inches(5.35), y, Inches(7.15),
                                Inches(1.05), text, fill=fill,
                                text_color=tcol, size=20, bold=True,
                                corner_pct=0.10, shadow=True)
        y = int(y + Inches(1.45))
    _add_arrow(slide, (Inches(8.9), Inches(3.12)),
               (Inches(8.9), Inches(3.48)), color=NAVY, weight_pt=2.5,
               head=True)
    _add_arrow(slide, (Inches(8.9), Inches(4.57)),
               (Inches(8.9), Inches(4.93)), color=NAVY, weight_pt=2.5,
               head=True)
    _draw_footer(slide, FOOTER_TEXT, 31)
    _add_pollbreak_badge(slide)
    return slide


# --------------------------------------------------------------------------
# Slide 34 — Solution (water)
# --------------------------------------------------------------------------

def slide_34_water_solution(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_OWN)
    _draw_action_title(slide, "Solution")
    # reworded and lifted by hand 2026-08-25 (was at y 1.85)
    _add_hierarchical_bullets(
        slide, Inches(0.42), Inches(1.48), Inches(12.4),
        Inches(1.33),
        [
            ("We know:", 0, {'bold': True}),
            ("Demand elasticity is −0.4", 1),
            ("Target quantity reduction = - 10%", 1),
        ],
        size=26, sub_size=24)
    eqs = [
        (_oED_frac('P'), NAVY),
        (_omml_text('−0.4') + _omml_text(' = ')
         + _omml_frac(_omml_text('−10%'), _o_pct('P')), NAVY),
        # 2026-08-26 (Nico): the FINAL solution line is dark red on
        # every solution slide; the steps above it stay navy
        (_o_pct('P') + _omml_text(' = ')
         + _omml_frac(_omml_text('−10%'), _omml_text('−0.4'))
         + _omml_text(' = 25% price increase'), RED),
    ]
    # equations and takeaway lifted 0.58" by hand 2026-08-25
    y = Inches(2.87)
    for omml, color in eqs:
        _add_math_equation(slide, Inches(3.25), y, Inches(6.8),
                           Inches(0.95), omml, size_pt=26, color=color)
        y = int(y + Inches(1.00))
    _add_takeaway_bar(slide, " Raise the price 25% to cut water use by 10%",
                      top=Inches(6.02), left=Inches(3.09),
                      width=Inches(6.96), height=Inches(0.66),
                      fill=GOLD, text_color=NAVY, size=18, bold=True,
                      rounded=True, shadow=True,
                      wingding_lead="\uf0e0")
    _draw_footer(slide, FOOTER_TEXT, 33)
    return slide


# --------------------------------------------------------------------------
# Slide 35 — Categories of own-price elasticity
# --------------------------------------------------------------------------

def slide_35_categories(prs):
    """Categories of own-price elasticity.

    2026-08-25: relaid out on CT's In-Class slide 28 - the definition on
    top, then three equal cards, each stating the category twice (once
    on the absolute value, once signed) and naming it.  The bars stay
    navy like the rest of the formula; what CT marks in gold is instead
    called out by the "absolute value" pointer.
    """
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_OWN)
    _draw_action_title(slide, "Categories of Own-Price Elasticity")
    _add_math_equation(
        slide, Inches(4.37), Inches(1.95), Inches(4.59), Inches(1.39),
        _oED_frac('P'), size_pt=28, color=NAVY, fill=CREAM, line=NAVY,
        rounded=True, shadow=True)

    cards = [
        (0.80, _omml_text('0 &lt; |') + _oED() + _omml_text('| &lt; 1'),
         _omml_text('\u22121 &lt; ') + _oED() + _omml_text(' &lt; 0'),
         "inelastic"),
        (4.82, _omml_text('|') + _oED() + _omml_text('| &gt; 1'),
         _oED() + _omml_text(' &lt; \u22121'), "elastic"),
        (8.84, _omml_text('|') + _oED() + _omml_text('| = 1'),
         _oED() + _omml_text(' = \u22121'), "unit elastic"),
    ]
    for x, top_omml, bot_omml, word in cards:
        _add_outlined_box(
            slide, Inches(x), Inches(3.98), Inches(3.69), Inches(1.86),
            "", line=NAVY, fill=WHITE, line_w=1.25, rounded=True,
            shadow=True, corner_pct=0.10)
        _add_math_equation(slide, Inches(x), Inches(4.08), Inches(3.69),
                           Inches(0.58), top_omml, size_pt=25,
                           color=NAVY)
        _add_math_equation(slide, Inches(x), Inches(4.66), Inches(3.69),
                           Inches(0.58), bot_omml, size_pt=25,
                           color=NAVY)
        _add_hierarchical_bullets(
            slide, Inches(x), Inches(5.28), Inches(3.69), Inches(0.44),
            [([("demand is ", {}),
               (word, {'bold': True, 'color': CBLUE})], 0,
              {'bullet_style': 'none', 'align': PP_ALIGN.CENTER})],
            size=23)

    # the pointer CT puts on the bars of the first card
    # nudged by hand 2026-08-25 (label was 1.30, arrow 2.42 -> 2.19)
    _add_text(slide, Inches(1.38), Inches(3.46), Inches(2.70),
              Inches(0.36), "absolute value", size=20, bold=True,
              italic=True, color=GOLD, font="Calibri",
              align=PP_ALIGN.CENTER)
    _add_arrow(slide, (Inches(2.59), Inches(3.85)),
               (Inches(2.36), Inches(4.23)), color=GOLD, weight_pt=2.0,
               head=True)
    _set_notes(slide, (
        "Three categories, and they are defined on the absolute value so "
        "the minus sign does not trip you up. Inelastic means the "
        "absolute value is below 1: quantity moves proportionally less "
        "than price, as with water at −0.4. Elastic means the "
        "absolute value is above 1: quantity moves proportionally more, "
        "as with the yoga lessons at about −6. Unit-elastic is the "
        "knife-edge case at exactly −1, where the two percentage "
        "changes offset each other. These labels matter because, as we "
        "will see, they decide whether a price increase raises or lowers "
        "revenue.\n\n"
        "Video link: https://www.youtube.com/watch?v=dmgFP0qteBU&t=32s "
        "(watch until around 3:00)"))
    _draw_footer(slide, FOOTER_TEXT, 34)
    return slide


# --------------------------------------------------------------------------
# Slide 36 — Example: yoga lessons (CorePower Yoga, CT item 7)
# --------------------------------------------------------------------------

def slide_36_yoga(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_OWN)
    _draw_action_title(slide, "Example: Demand Elasticity for Yoga Lessons")
    # text and photo both lifted by hand 2026-08-25
    _add_hierarchical_bullets(
        slide, Inches(0.45), Inches(1.65), Inches(12.4),
        Inches(1.9),
        [
            ("CorePower Yoga cuts its price by 17%", 0),
            ([("The number of booked lessons ", {}),
              ("double", {'bold': True}), (" as a result", {})], 0, {}),
            ("What is the price elasticity of demand that the yoga "
             "studio faces?", 0),
        ],
        size=26)
    _add_media_image(slide, "image32.jpeg", left=Inches(2.67),
                     top=Inches(3.65), width=Inches(8.0))
    _draw_footer(slide, FOOTER_TEXT, 35)
    _add_pollbreak_badge(slide)
    return slide


# --------------------------------------------------------------------------
# Slide 39 — Solution (yoga)
# --------------------------------------------------------------------------

def slide_39_yoga_solution(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_OWN)
    _draw_action_title(slide, "Solution")
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.15), Inches(1.85), Inches(12.4),
        Inches(0.6),
        # 2026-08-25 (Nico): the "%" is what makes it unit-free, so it
        # is bold on a yellow highlight
        [([("Unit free  →  ", {}), ("not", {'bold': True}),
           (" 6", {}),
           ("%", {'bold': True, 'highlight': "FFFF00"})], 0, {})],
        size=26)
    _add_math_equation(
        slide, Inches(1.65), Inches(2.70), Inches(10.0), Inches(1.25),
        # 2026-08-25 (Nico): capital D subscript.  2026-08-26: the
        # final solution is dark red on every solution slide.
        _oED() + _omml_text(' = ')
        + _omml_frac(_o_pct('Q'), _o_pct('P')) + _omml_text(' = ')
        + _omml_frac(_omml_text('+100%'), _omml_text('−17%'))
        + _omml_text('  =  −5.88   (or about −6)'),
        size_pt=28, color=RED)
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.15), Inches(4.35), Inches(12.4),
        Inches(2.2),
        [
            ("Is demand for yoga lessons elastic or inelastic?", 0),
            ([("Answer: ", {}),
              ("elastic", {'bold': True,
                           'color': RGBColor(0x00, 0x70, 0xC0)}),
              (", as ", {}), ("Eᴅ", {'italic': True}),
              (" < −1   (or |", {}), ("Eᴅ", {'italic': True}),
              ("| > 1)", {})], 1, {}),
            ("As the price falls, quantity rises by MORE than the price "
             "falls", 1),
        ],
        size=26, sub_size=24)
    _draw_footer(slide, FOOTER_TEXT, 37)
    return slide


# --------------------------------------------------------------------------
# Slide 40 — Two ways to compute elasticity: Method 1
# --------------------------------------------------------------------------

def slide_40_method1(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_OWN)
    # retitled by hand 2026-08-25
    _custom_title_runs(slide, [
        ("Computing Elasticity:  ", {}),
        ("Method 1", {'color': RED})])
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.15), Inches(1.85), Inches(12.4),
        Inches(0.65),
        # 2026-08-25 (Nico): real subscripts, not a smaller digit -
        # apply_subscripts() turns the ₀/₁ characters into true
        # baseline-shifted runs deck-wide
        [([("When you have two price/quantity observations:  (", {}),
           ("P₀", {'italic': True}), (", ", {}),
           ("Q₀", {'italic': True}), (")  and  (", {}),
           ("P₁", {'italic': True}), (", ", {}),
           ("Q₁", {'italic': True}), (")", {})],
          0, {})],
        size=24)
    # equation narrowed and moved left by hand 2026-08-25 (was
    # 3.35 / 2.75, 6.6 x 1.95) to make room for the convention card
    _add_math_equation(
        slide, Inches(2.48), Inches(2.58), Inches(4.68), Inches(2.14),
        _oED_frac('P') + _omml_text(' = ')
        + _omml_frac(
            _omml_frac(_omml_sub(_omml_run('Q'), _omml_text('1'))
                       + _omml_text('−')
                       + _omml_sub(_omml_run('Q'), _omml_text('0')),
                       _omml_sub(_omml_run('Q'), _omml_text('0'))),
            _omml_frac(_omml_sub(_omml_run('P'), _omml_text('1'))
                       + _omml_text('−')
                       + _omml_sub(_omml_run('P'), _omml_text('0')),
                       _omml_sub(_omml_run('P'), _omml_text('0')))),
        size_pt=28, color=NAVY, fill=CREAM, line=NAVY, rounded=True,
        shadow=True)
    # 2026-08-25 (Nico): this convention is easy to miss and costs
    # students the sign, so it gets its own louder card - concept blue
    # on a pale blue fill, a heavier border and a bigger label than the
    # ordinary cream callouts
    # card moved beside the formula by hand 2026-08-25 (was at
    # 8.30 / 4.86, 4.55 x 1.16)
    _add_rounded_filled_box(
        slide, Inches(8.12), Inches(2.77), Inches(2.96), Inches(1.16),
        "", fill=PALE_BLUE, line=CBLUE, line_w=2.25, corner_pct=0.14,
        shadow=True)
    _add_hierarchical_bullets(
        slide, Inches(8.24), Inches(2.88), Inches(2.73), Inches(0.95),
        [("CONVENTION", 0,
          {'bullet_style': 'none', 'align': PP_ALIGN.CENTER,
           'size': 14, 'bold': True, 'color': CBLUE}),
         ([("always relative to the ", {}),
           ("initial", {'bold': True, 'underline': True}),
           (" point  (P₀, Q₀)", {})], 0,
          {'bullet_style': 'none', 'align': PP_ALIGN.CENTER,
           'size': 17, 'color': NAVY, 'space_before_pts': 2})],
        size=17)
    # 2026-08-28 (Nico): the caveat leaves the bullet list and becomes a
    # navy box, so the limitation of Method 1 reads as a statement in
    # its own right.  Its sibling sits on slide 42 (Method 2, exact) -
    # same geometry and type size, so the two read as a matched pair.
    # The one word that carries the point is set in gold.
    _add_convention_box(
        slide, MARGIN + Inches(0.15), Inches(5.15), Inches(8.50),
        Inches(0.75),
        runs=[("This method ", {'color': WHITE}),
              ("approximates", {'bold': True, 'color': GOLD}),
              (" the percentage changes", {'color': WHITE})],
        size=22, align=PP_ALIGN.CENTER, corner_pct=0.14,
        fill_rgb=NAVY, border=NAVY)
    # 2026-08-25 (Nico): the TA video pointer is a reference box now
    _add_reference_box(
        slide, Inches(4.30), Inches(6.36), Inches(3.90), Inches(0.62),
        "TA Math Review Videos", kind="video", size=16)
    # 2026-08-24: the deck-wide reference-box convention (glyph + the
    # problem-set number only).  2026-08-26: problem-set pointers live
    # in the bottom-RIGHT corner deck-wide (PS_BOX_XY)
    _add_reference_box(slide, PS_BOX_XY[0], PS_BOX_XY[1], Inches(3.0),
                       Inches(0.5), "Problem Set 2", kind="ps", size=18)
    _draw_footer(slide, FOOTER_TEXT, 38)
    # 2026-08-25 (Nico): two arrows from the convention card to the
    # two denominators it is talking about
    _add_arrow(slide, (Inches(8.026), Inches(3.070)),
               (Inches(6.477), Inches(3.442)), color=CBLUE,
               weight_pt=2.0, head=True)
    _add_arrow(slide, (Inches(8.026), Inches(3.721)),
               (Inches(6.477), Inches(4.384)), color=CBLUE,
               weight_pt=2.0, head=True)
    return slide


# --------------------------------------------------------------------------
# Slide 41 — Amazon e-books (poll setup; keeps its live poll)
# --------------------------------------------------------------------------

def slide_41_ebooks(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_OWN)
    _draw_action_title(slide, "Amazon E-Books")
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.15), Inches(2.30), Inches(10.3),
        Inches(4.1),
        [
            ("Amazon said that its internal data showed that when an "
             "e-book is priced at $9.99, it sells nearly twice as many "
             "copies as when it is priced at $14.99. It argued that, as "
             "a result, total revenue at $9.99 is more than when the "
             "book is priced higher. “At $9.99, the total pie is "
             "bigger,” stated Amazon.", 0,
             {'italic': True, 'bullet_style': 'none'}),
            ("", 0, {'bullet_style': 'none'}),
            ("What is the implied price elasticity of demand (using "
             "$14.99 as the initial price)?", 0),
        ],
        size=24, line_spacing_pts=16)
    _add_media_image(slide, "image37.png", left=Inches(11.15),
                     top=Inches(1.55), width=Inches(1.6),
                     rounded=False, shadow=False)
    _set_notes(slide, (
        "This is an elasticity argument made in public, during Amazon's "
        "2014 pricing dispute with the publisher Hachette. Amazon claimed "
        "its internal data showed an e-book priced at $9.99 sells nearly "
        "twice as many copies as the same book at $14.99, and concluded "
        "that total revenue is higher at the lower price – “at "
        "$9.99, the total pie is bigger.” That is a claim about "
        "elastic demand, and you can now check it. Compute the implied "
        "elasticity using $14.99 as the initial price.\n\n"
        "https://www.wsj.com/articles/amazon-calls-for-hachette-to-cut-"
        "e-book-prices-1406675179"))
    _draw_footer(slide, FOOTER_TEXT, 43)
    _add_pollbreak_badge(slide)
    return slide


# --------------------------------------------------------------------------
# Slide 44 — Solution (e-books)
# --------------------------------------------------------------------------

def slide_44_ebooks_solution(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_OWN)
    _draw_action_title(slide, "Solution")
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.15), Inches(1.85), Inches(12.4),
        Inches(1.35),
        [
            ([("Problem states: Go from 14.99 to 9.99  →  ", {}),
              ("100% more sales", {'bold': True})], 0, {}),
            ([("Thus: Use 14.99 as initial price (", {}),
              ("P", {'italic': True}), ("0", {'size': 16}),
              (")  and  ", {}), ("P", {'italic': True}),
              ("1", {'size': 16}), (" = 9.99", {})], 0, {}),
        ],
        size=24)
    _add_math_equation(
        slide, Inches(2.15), Inches(3.45), Inches(9.0), Inches(1.85),
        _oED() + _omml_text(' = ')
        + _omml_frac(
            _omml_frac(_omml_sub(_omml_run('Q'), _omml_text('1'))
                       + _omml_text('−')
                       + _omml_sub(_omml_run('Q'), _omml_text('0')),
                       _omml_sub(_omml_run('Q'), _omml_text('0'))),
            _omml_frac(_omml_sub(_omml_run('P'), _omml_text('1'))
                       + _omml_text('−')
                       + _omml_sub(_omml_run('P'), _omml_text('0')),
                       _omml_sub(_omml_run('P'), _omml_text('0'))))
        + _omml_text(' = ')
        + _omml_frac(_omml_text('100%'),
                     _omml_frac(_omml_text('9.99 − 14.99'),
                                _omml_text('14.99')))
        + _omml_text(' = ')
        + _omml_frac(_omml_text('100%'), _omml_text('−33%'))
        + _omml_text(' = −3'),
        size_pt=24, color=NAVY)
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.15), Inches(5.75), Inches(12.4),
        Inches(0.6),
        [([("→  ", {}),
           ("E", {'italic': True, 'color': RED, 'bold': True}),
           ("D", {'size': 16, 'color': RED, 'bold': True}),
           (" ≈ −3 :  demand for e-books is elastic",
            {'color': RED, 'bold': True})], 0,
          {'bullet_style': 'none'})],
        size=26)
    _set_notes(slide, (
        "Going from $14.99 to $9.99 is a price change of about "
        "−33%, and the quantity roughly doubles, so +100%. Divide "
        "and you get an elasticity of about −3. That is well below "
        "−1, so demand for these e-books is elastic, which is "
        "exactly what Amazon's “bigger pie” claim requires: "
        "with elastic demand a price cut raises total revenue. Worth "
        "noting where the number came from – it rests on Amazon's "
        "own internal sales data as reported at the time, not on an "
        "independent estimate.\n\n"
        "Rounding to the closest integer"))
    _draw_footer(slide, FOOTER_TEXT, 46)
    return slide


# --------------------------------------------------------------------------
# Slide 45 — Mega Millions example (NEW; corrected facts, decision 2026-08-14)
# --------------------------------------------------------------------------

def slide_45_megamillions(prs):
    """Lottery Sales - the setup (2026-08-25, on CT's In-Class p.15).

    CT dates the change to 2024 and credits "the MA State Lottery".
    Both are off: the $2 -> $5 price was a NATIONAL Mega Millions
    redesign that took effect on 8 April 2025 (announced in October
    2024); the state lotteries sell the tickets but did not set the
    price.  The numbers here are New York's, from the Hansen / Misra /
    Singh working paper, and give a very different elasticity from CT's
    -0.2 - see the note on the solution slide.
    """
    slide = _blank_slide(prs)
    # 2026-08-25 (Nico): CT's full-bleed lottery photo as a washed-out
    # background, drawn FIRST so everything else sits on top.  The wash
    # (alphaModFix) is what keeps the navy text readable.
    _add_media_image(slide, "ct_lottery_bg.png", left=0, top=0,
                     width=SLIDE_W, rounded=False, shadow=False,
                     transparency=78)
    _draw_top_bar_tc(slide, TAG_OWN)
    _draw_action_title(slide, "Lottery Sales: Mega Millions Raises Its Price")
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.15), Inches(2.05), Inches(12.4),
        Inches(3.0),
        [
            ([("April 2025: Mega Millions raises its ticket price from "
               "$2 to $5", {})], 0, {}),
            ([("In New York, tickets sold per drawing fall from about "
               "1.9 million to about 560,000", {})], 0, {}),
            ([("What is the implied price elasticity of demand ",
               {'bold': True}),
              ("(using $2 as the initial price)?", {})], 0, {}),
        ],
        size=26, line_spacing_pts=20)
    _add_text(slide, MARGIN, Inches(6.78), Inches(9.0), Inches(0.3),
              "Source: Hansen, Misra & Singh, \u201cPricing a "
              "Participation-Dependent Product: Evidence from the Mega "
              "Millions Redesign\u201d", size=12, italic=True,
              color=GRAY, font="Calibri")
    _set_notes(slide, (
        "A second real-world check, and this one is recent. On 8 April "
        "2025 Mega Millions raised its ticket price from $2 to $5 - a "
        "150% increase, and only the second price change in the game's "
        "history. Hansen, Misra and Singh collected the sales data: in "
        "New York, tickets sold per drawing fell from about 1.9 million "
        "to about 560,000. Ask them to work out the implied elasticity "
        "before showing the solution, and remind them to measure both "
        "changes from the initial point."))
    _draw_footer(slide, FOOTER_TEXT, 39)
    _add_pollbreak_badge(slide)
    return slide


def slide_41_megamillions_solution(prs):
    """The worked answer, laid out like CT's In-Class p.18 but in native
    OMML and in the deck's navy rather than red."""
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_OWN)
    _draw_action_title(slide, "Solution: Mega Millions")
    _add_hierarchical_bullets(
        slide, Inches(1.00), Inches(1.55), Inches(11.33), Inches(1.05),
        [
            ([("Price:  ", {'bold': True}),
              ("P\u2080 = $2   \u2192   P\u2081 = $5", {})], 0,
             {'bullet_style': 'none'}),
            ([("Quantity:  ", {'bold': True}),
              ("Q\u2080 = 1,900,000   \u2192   Q\u2081 = 560,000", {})], 0,
             {'bullet_style': 'none'}),
        ],
        size=22, line_spacing_pts=8)

    _add_math_equation(
        slide, Inches(1.15), Inches(2.86), Inches(11.0), Inches(1.05),
        _o_pct('Q') + _omml_text(' = ')
        + _omml_frac(_omml_sub(_omml_run('Q'), _omml_text('1'))
                     + _omml_text(' \u2212 ')
                     + _omml_sub(_omml_run('Q'), _omml_text('0')),
                     _omml_sub(_omml_run('Q'), _omml_text('0')))
        + _omml_text(' = ')
        + _omml_frac(_omml_text('560,000 \u2212 1,900,000'),
                     _omml_text('1,900,000'))
        + _omml_text(' = \u221270.5%'),
        size_pt=21, color=NAVY)
    _add_math_equation(
        slide, Inches(1.15), Inches(4.06), Inches(11.0), Inches(1.05),
        _o_pct('P') + _omml_text(' = ')
        + _omml_frac(_omml_sub(_omml_run('P'), _omml_text('1'))
                     + _omml_text(' \u2212 ')
                     + _omml_sub(_omml_run('P'), _omml_text('0')),
                     _omml_sub(_omml_run('P'), _omml_text('0')))
        + _omml_text(' = ')
        + _omml_frac(_omml_text('5 \u2212 2'), _omml_text('2'))
        + _omml_text(' = +150%'),
        size_pt=21, color=NAVY)
    _add_math_equation(
        slide, Inches(1.15), Inches(5.26), Inches(11.0), Inches(1.05),
        _oED() + _omml_text(' = ')
        + _omml_frac(_o_pct('Q'), _o_pct('P')) + _omml_text(' = ')
        + _omml_frac(_omml_text('\u221270.5%'), _omml_text('150%'))
        + _omml_text('  \u2248  \u22120.47'),
        # 2026-08-26 (Nico): final solution in dark red
        size_pt=21, color=RED)

    _add_convention_box(
        slide, Inches(1.85), Inches(6.42), Inches(9.6), Inches(0.62),
        prefix="Caution: ",
        body="Method 1 approximates % changes \u2013 best for small price "
             "changes. Here the changes are large, so \u22120.47 is a rough "
             "approximation", size=15)
    _set_notes(slide, (
        "Both percentage changes are measured from the initial point, "
        "as always with Method 1. Quantity falls by about 70%, price "
        "rises by 150%, so the implied elasticity is about \u22120.47 - "
        "well inside the inelastic range, which is why revenue still "
        "rose even though the number of tickets collapsed.\n\n"
        "If you compare this with CT's version of the example: she "
        "dates the increase to 2024 and attributes it to the "
        "Massachusetts State Lottery, and her sales figures give an "
        "elasticity of about \u22120.2. The price change was in fact a "
        "national Mega Millions redesign effective 8 April 2025, and "
        "the New York figures above come from the Hansen / Misra / "
        "Singh study."))
    _draw_footer(slide, FOOTER_TEXT, 43)
    return slide




# ==========================================================================
#  BATCH D — Method 2 / point elasticity / linear case / Uber /
#             special cases / determinants / market-vs-firm (46–58)
# ==========================================================================

def _oD_slope():
    return _omml_frac(_omml_text('Δ') + _omml_run('Q'),
                      _omml_text('Δ') + _omml_run('P'))


def _oD_deriv():
    return _omml_frac(_omml_run('dQ'), _omml_run('dP'))


# --------------------------------------------------------------------------
# Slide 46 — Two ways to compute elasticity: Method 2 ("Point Elasticity")
# --------------------------------------------------------------------------

def slide_46_method2(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_OWN)
    # retitled by hand 2026-08-25
    _custom_title_runs(slide, [
        ("Computing Elasticity:  ", {}),
        ("Method 2", {'color': RED}),
        ("  (“Point Elasticity”)", {'size': 24})])
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.15), Inches(1.80), Inches(12.4),
        Inches(0.6),
        [("When you know the entire demand curve, use it!", 0,
          {'bold': True})],
        size=24)
    _add_math_equation(
        slide, Inches(3.43), Inches(2.39), Inches(5.70), Inches(2.10),
        _oED_frac('P') + _omml_text(' = ')
        + _omml_frac(_omml_frac(_omml_text('Δ') + _omml_run('Q'),
                                _omml_run('Q')),
                     _omml_frac(_omml_text('Δ') + _omml_run('P'),
                                _omml_run('P')))
        + _omml_text(' = ') + _oD_slope() + _omml_text(' ∙ ')
        + _omml_frac(_omml_run('P'), _omml_run('Q')),
        size_pt=28, color=NAVY, fill=CREAM, line=NAVY, rounded=True,
        shadow=True)
    # 2026-08-25 (Nico, on CT's In-Class slide 38): ring the
    # ΔQ/ΔP term and name it, in dark red rather than CT's gold
    _add_oval_outline(slide, Inches(7.42), Inches(2.96), Inches(0.80),
                      Inches(1.06), color=DARK_RED, weight_pt=2.25,
                      shadow=False)
    _add_hierarchical_bullets(
        slide, Inches(9.55), Inches(3.06), Inches(3.35), Inches(0.85),
        [("Slope of the demand curve", 0,
          {'bullet_style': 'none', 'size': 17, 'bold': True,
           'italic': True, 'color': DARK_RED}),
         ("(constant for a linear demand curve)", 0,
          {'bullet_style': 'none', 'size': 15, 'italic': True,
           'color': DARK_RED, 'space_before_pts': 0})],
        size=17)
    _add_arrow(slide, (Inches(9.48), Inches(3.36)),
               (Inches(8.30), Inches(3.44)), color=DARK_RED,
               weight_pt=1.75, head=True)
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.15), Inches(4.65), Inches(12.4),
        Inches(2.1),
        [
            # 2026-08-26 (Nico): the ΔQ/ΔP term is set in red so it ties
            # back to the red ring around it in the formula above
            ([("ΔQ/ΔP", {'color': RED}),
              (" is the slope of the demand curve", {})], 0, {}),
            ([("Mathematically: Derivative of ", {}),
              ("Q", {'italic': True}), (" with respect to ", {}),
              ("P", {'italic': True}), (" :  ", {}),
              ("dQ/dP", {'italic': True})], 1, {}),
            ("See TA Math Review videos", 1),
            # 2026-08-28 (Nico's hand edit): the fourth bullet ("For a
            # linear demand curve, the slope is constant") is deleted -
            # the ringed callout beside the formula already says it.
            # The box below replaces it.
        ],
        size=24, sub_size=22, line_spacing_pts=10)
    # 2026-08-28 (Nico): the counterpart of slide 38's navy box - same
    # geometry, same type size, so Method 1 (approximate) and Method 2
    # (exact) read as a matched pair
    _add_convention_box(
        slide, MARGIN + Inches(0.15), Inches(6.10), Inches(8.50),
        Inches(0.75),
        runs=[("This method computes the ", {'color': WHITE}),
              ("exact", {'bold': True, 'color': GOLD}),
              (" elasticity at a given (", {'color': WHITE}),
              ("P", {'italic': True, 'color': WHITE}),
              (", ", {'color': WHITE}),
              ("Q", {'italic': True, 'color': WHITE}),
              (") point", {'color': WHITE})],
        size=22, align=PP_ALIGN.CENTER, corner_pct=0.14,
        fill_rgb=NAVY, border=NAVY)
    # 2026-08-26 (Nico): the pointer moves to the bottom-right corner,
    # the deck's default position for a post-work reference box
    _add_reference_box(slide, PS_BOX_XY[0], PS_BOX_XY[1], Inches(3.00),
                       Inches(0.50), "Problem Set 2", kind="ps", size=17)
    _draw_footer(slide, FOOTER_TEXT, 42)
    return slide


# --------------------------------------------------------------------------
# Slide 47 — Computing the point elasticity: step-by-step
# --------------------------------------------------------------------------

def slide_47_point_steps(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_OWN)
    # 2026-08-26 (Nico): the title now names the method, with "Method 2"
    # in red exactly as on slide 42
    _custom_title_runs(slide, [
        ("Computing the Point Elasticity: ", {}),
        ("Method 2", {'color': RED}),
        (" – Step-by-Step", {})], size=32)
    # given: inverse demand + task (top-left cream card).  2026-08-26:
    # wider card, text up from 17 to 24 pt
    _add_convention_box(
        slide, MARGIN + Inches(0.15), Inches(1.70), Inches(6.644),
        Inches(1.15),
        # 2026-08-25 (Nico): start from the DEMAND function itself,
        # so there is no "solve for Q" step to do first
        runs=[("Demand function:   Q = 400 − 4 P", {'bold': True}),
              ("\nTask: Compute the point elasticity at P = 25",
               {'bold': True})],
        size=24)
    steps = [
        ("Step 1:  Compute derivative dQ/dP",
         _oD_deriv() + _omml_text(' = −4')),
        ("Step 2:  Obtain Q at P = 25",
         _omml_run('Q') + _omml_text(' = 400 − 4 × 25 = 300')),
        ("Step 3:  Put everything into the formula",
         _oED() + _omml_text(' = ') + _oD_deriv() + _omml_text(' ∙ ')
         + _omml_frac(_omml_run('P'), _omml_run('Q'))
         + _omml_text(' = −4 × ')
         + _omml_frac(_omml_text('25'), _omml_text('300'))
         + _omml_text(' = −')
         + _omml_frac(_omml_text('1'), _omml_text('3'))),
    ]
    # three steps now, so they sit lower and breathe more
    y = Inches(3.45)
    for label, omml in steps:
        _add_hierarchical_bullets(
            slide, MARGIN + Inches(0.15), int(y + Inches(0.12)),
            Inches(5.9), Inches(0.7),
            [(label, 0, {'bold': True, 'bullet_style': 'none'})],
            size=20)
        _add_math_equation(slide, Inches(6.55), y, Inches(5.9),
                           Inches(0.85), omml, size_pt=20, color=NAVY,
                           fill=CREAM, line=NAVY, rounded=True)
        y = int(y + Inches(1.00))
    _draw_footer(slide, FOOTER_TEXT, 43)
    return slide


# --------------------------------------------------------------------------
# Slide 48 — Obtain price elasticity from the demand function (poll setup)
# --------------------------------------------------------------------------

def slide_48_from_demand_fn(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_OWN)
    _draw_action_title(
        slide, "Obtain Price Elasticity from the Demand Function")
    # 2026-08-26 (Nico): three lines on an otherwise empty slide, so they
    # go up to 28 pt and get 18 pt of air before each — the sparse-slide
    # spacing rule in the Teaching CLAUDE.md
    _add_mixed_textbox(
        slide, MARGIN + Inches(0.15), Inches(2.30), Inches(12.4),
        Inches(2.6),
        [
            ("text", "A firm’s demand function is  ", {'size': 28}),
            ("omml", _omml_run('Q') + _omml_text(' = 10 − ')
             + _omml_run('P'), {'size': 28}),
            ("break", None, None),
            ("text", "The firm charges a price of  ", {'size': 28}),
            ("omml", _omml_run('P') + _omml_text(' = 2'), {'size': 28}),
            ("break", None, None),
            ("text", "What is its elasticity of demand at this point?",
             {'size': 28, 'bold': True}),
        ],
        default_size=28, space_before_pts=18)
    _draw_footer(slide, FOOTER_TEXT, 44)
    _add_pollbreak_badge(slide)
    return slide


# --------------------------------------------------------------------------
# Slide 51 — Solution: Q = 10 − P at P = 2
# --------------------------------------------------------------------------

def slide_51_qp_solution(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_OWN)
    _custom_title_runs(slide, [
        ("Solution: If  ", {}),
        ("Q = 10 − P", {'italic': True}),
        (" , Elasticity at  ", {}),
        ("P = 2", {'italic': True}), (" ?", {})])
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.15), Inches(1.90), Inches(12.4),
        Inches(2.5),
        [
            # 2026-08-25 (Nico): the steps renumber with slide 43 -
            # there is no "solve for Q" step any more
            ([("Answer is ", {}), ("−0.25", {'bold': True})], 0, {}),
            ([("Step 1:  ", {'bold': True}), ("Q", {'italic': True}),
              (" = 10 − ", {}), ("P", {'italic': True}),
              ("   and thus   ", {}), ("dQ/dP", {'italic': True}),
              (" = −1", {})], 0, {}),
            ([("Step 2:  ", {'bold': True}), ("Compute ", {}),
              ("Q", {'italic': True}), (" at ", {}),
              ("P", {'italic': True}), (" = 2 :   ", {}),
              ("Q", {'italic': True}), (" = 10 − 2 = 8", {})], 0, {}),
            ("Step 3:", 0, {'bold': True}),
        ],
        size=24, line_spacing_pts=12)
    # 2026-08-26 (Nico): the formula sits directly under "Step 3:"
    _add_math_equation(
        slide, Inches(2.70), Inches(3.64), Inches(7.6), Inches(1.15),
        _oED() + _omml_text(' = ') + _oD_deriv() + _omml_text(' ∙ ')
        + _omml_frac(_omml_run('P'), _omml_run('Q'))
        + _omml_text(' = −1 × ')
        + _omml_frac(_omml_text('2'), _omml_text('8'))
        + _omml_text(' = −')
        + _omml_frac(_omml_text('1'), _omml_text('4')),
        # 2026-08-26 (Nico): final solution in dark red
        size_pt=26, color=RED)
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.15), Inches(6.05), Inches(12.4),
        Inches(0.55),
        [([("→  Demand is ", {}),
           ("inelastic", {'bold': True,
                          'color': RGBColor(0x00, 0x70, 0xC0)})], 0,
          {'bullet_style': 'none'})],
        size=26)
    _draw_footer(slide, FOOTER_TEXT, 47)
    return slide


# --------------------------------------------------------------------------
# Slide 52 — Elasticity along the demand curve: linear case
# --------------------------------------------------------------------------

def slide_52_linear_elasticity(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_OWN)
    _draw_action_title(
        slide, "Elasticity Along the Demand Curve: Linear Case")
    _add_math_equation(
        slide, Inches(8.45), Inches(1.60), Inches(4.3), Inches(1.15),
        _oED_frac('P') + _omml_text(' = ') + _oD_slope()
        + _omml_text(' ∙ ') + _omml_frac(_omml_run('P'), _omml_run('Q')),
        size_pt=24, color=NAVY, fill=CREAM, line=NAVY, rounded=True,
        shadow=True)
    # 2026-08-25 (Nico): the graph is now CT's In-Class slide 39,
    # element for element.  Her coordinates are kept verbatim below and
    # shifted by CT_DX / CT_DY, which slides the whole graph left and up
    # so it clears our right-hand column (she has no column - her
    # elasticity formula sits inside the plot area).
    CT_DX, CT_DY = -1.75, -0.60

    def CX(v):
        return Inches(v + CT_DX)

    def CY(v):
        return Inches(v + CT_DY)

    # axes (navy, triangle heads) with her word labels
    _add_arrow(slide, (CX(2.550), CY(6.550)), (CX(2.550), CY(3.250)),
               color=NAVY, weight_pt=1.8, head=True)
    _add_arrow(slide, (CX(2.550), CY(6.550)), (CX(9.250), CY(6.550)),
               color=NAVY, weight_pt=1.8, head=True)
    _add_text(slide, CX(2.200), CY(2.660), Inches(1.40), Inches(0.30),
              "Price", size=19.2, bold=True, italic=True, color=NAVY,
              font="Calibri")
    _add_text(slide, CX(8.550), CY(6.620), Inches(1.60), Inches(0.30),
              "Quantity", size=19.2, bold=True, italic=True, color=NAVY,
              font="Calibri")
    # the demand curve
    _add_arrow(slide, (CX(2.590), CY(3.450)), (CX(8.346), CY(6.550)),
               color=NAVY, weight_pt=2.6, head=False)
    _add_text(slide, CX(7.416), CY(6.194), Inches(0.50), Inches(0.32),
              "D", size=21.6, bold=True, italic=True, color=NAVY,
              font="Calibri")
    # the two braces, one along each stretch, and their region labels
    _add_rot_brace(slide, CX(3.950), CY(2.380), Inches(0.42),
                   Inches(3.20), 298.0, color=GOLD, weight_pt=2.0)
    _add_rot_brace(slide, CX(6.830), CY(3.930), Inches(0.42),
                   Inches(3.20), 298.0, color=GOLD, weight_pt=2.0)
    _add_runs_text(
        slide, CX(4.144), CY(3.340), Inches(3.20), Inches(0.40),
        [("Eᴅ < −1   ", {}), ("elastic", {'bold': True, 'color': CBLUE})],
        size=21.6, color=NAVY)
    _add_runs_text(
        slide, CX(7.253), CY(4.954), Inches(3.40), Inches(0.40),
        [("−1 < Eᴅ < 0   ", {}),
         ("inelastic", {'bold': True, 'color': CBLUE})],
        size=21.6, color=NAVY)
    # the midpoint, called out in gold with a leader to the dot
    _add_runs_text(
        slide, CX(6.156), CY(4.280), Inches(3.20), Inches(0.34),
        [("Eᴅ = −1   unit elastic", {})],
        size=20.4, bold=True, color=GOLD)
    _add_arrow(slide, (CX(6.050), CY(4.620)), (CX(5.550), CY(4.960)),
               color=GOLD, weight_pt=1.4, head=False)
    _add_oval_filled(slide, CX(5.370), CY(4.900), Inches(0.20),
                     Inches(0.20), fill=GOLD, line=NAVY, weight_pt=1.0)
    # the two intercepts, each with a short leader onto the curve
    _add_runs_text(
        slide, CX(2.750), CY(2.887), Inches(1.70), Inches(0.30),
        [("Eᴅ = −∞", {})], size=18, color=NAVY)
    _add_arrow(slide, (CX(2.745), CY(3.151)), (CX(2.600), CY(3.429)),
               color=NAVY, weight_pt=1.5, head=True)
    _add_runs_text(
        slide, CX(8.755), CY(6.118), Inches(1.00), Inches(0.30),
        [("Eᴅ = 0", {})], size=18, color=NAVY)
    _add_arrow(slide, (CX(8.637), CY(6.340)), (CX(8.377), CY(6.500)),
               color=NAVY, weight_pt=1.5, head=True)
    # 2026-08-26 (Nico): the label heads its own underlined line, the
    # statement follows underneath
    _add_hierarchical_bullets(
        slide, Inches(8.525), Inches(2.915), Inches(4.3), Inches(1.0),
        [([("Linear demand curve:", {'underline': True})], 0,
          {'bullet_style': 'none'}),
         ([("E", {'italic': True}), ("D", {'size': 16}),
           (" goes from −∞ (approaching y-axis) to 0 (approaching "
            "x-axis)", {})], 0,
          {'bullet_style': 'none', 'space_before_pts': 6})],
        size=18)
    # 2026-08-26 (Nico): this line is a WARNING, not a convention —
    # transparent dark-red wash, dark-red border, dark-red pointer
    _add_convention_box(
        slide, Inches(8.518), Inches(4.90), Inches(4.3), Inches(0.80),
        runs=[("We will learn: ", {}),
              ("Firms should NOT operate", {'bold': True}),
              (" in the inelastic area (Module 2, Videos 1+2)", {})],
        size=15, fill_rgb=RED, fill_alpha_pct=12, border=RED)
    _add_arrow(slide, (Inches(8.518), Inches(5.242)),
               (Inches(7.718), Inches(4.796)), color=RED, weight_pt=2.0,
               head=True)
    _set_notes(slide, (
        "Unless the demand curve has the shape of a hyperbola, the "
        "own-price elasticity will be different at different points. "
        "Recall that the own-price elasticity of demand is E_d = %ΔQ/%ΔP "
        "= ΔQ/ΔP ∙ P/Q. The first term is the inverse slope of the "
        "demand curve. Since the slope is constant in a linear demand "
        "function, the second term in the elasticity formula, P/Q, will "
        "determine how the elasticity varies along the curve. We can "
        "decompose the curve into three parts. Elastic: In the upper "
        "left portion of the demand curve, price is very high relative "
        "to quantity, making P/Q large. In the elasticity formula, the "
        "inverse slope (which is negative) is multiplied by a large "
        "number, and the product will be less than negative one. For "
        "this reason, the upper left portion of the demand curve is "
        "elastic. Inelastic: In the lower right portion of the demand "
        "curve, price is low relative to quantity, making P/Q small. In "
        "the elasticity formula, the inverse slope (which is negative) "
        "is multiplied by a small number, and the product will be "
        "greater than negative one. For this reason, the lower right "
        "portion of the demand curve is inelastic. Unit elastic: Unit "
        "elasticity is the point on the demand curve at which the "
        "elasticity of demand is equal to −1. At this point, the "
        "percentage change in price is perfectly offset by the "
        "percentage change in quantity. In addition to these three "
        "regions, we can also compute the elasticity where the demand "
        "curve intersects the x and y axes. At the intersection with "
        "the y-axis, the quantity is zero, so P/Q = ∞, and the "
        "elasticity is negative infinity. The intuition behind this is "
        "that a small decrease in price will yield a change in quantity "
        "demanded from zero to something positive. Even if this change "
        "in quantity is small, the percentage change is infinite "
        "because it is rising from zero. At the intersection with the "
        "x-axis, the price is zero, so P/Q = 0, and the elasticity of "
        "demand is also zero."))
    _draw_footer(slide, FOOTER_TEXT, 48)
    return slide


# --------------------------------------------------------------------------
# Slide 53 — Important insight (don't operate in the inelastic part)
# --------------------------------------------------------------------------

def slide_53_insight(prs):
    slide = make_content_bulleted(
        prs, 49, TAG_OWN,
        "Important Insight",
        [
            ([("Intuition:", {'color': RED, 'bold': True}),
              ("  With inelastic demand, quantity reacts little when "
               "prices are raised  →  higher prices will ", {}),
              ("raise", {'underline': True}), (" revenues", {})], 0, {}),
            ("Also, quantity declines, so (variable) costs decline", 1,
             {'size': 20}),
            ("Thus: Firms can increase profits by raising prices", 1,
             {'size': 20}),
            ("Example: A study in 2016 estimated Uber’s price "
             "elasticity as −0.4 (inelastic)", 0),
            ("Subsequently, Uber raised its prices by an average of 18% "
             "per year in the 2018-2022 period. Much more than average "
             "inflation over the same period (4.5%)", 1, {'size': 20}),
            ("Now prices (and revenues) are much higher", 1,
             {'size': 20}),
        ],
        size=24, sub_size=22, bullets_width=Inches(8.60),
        bullets_top=Inches(2.25))
    # 2026-08-26 (Nico): the headline claim leaves the bullet list and
    # becomes a slide-wide WARNING box — same transparent dark red as
    # the "We will learn" box on slide 48 — with the bullets below it
    _add_convention_box(
        slide, MARGIN, Inches(1.45), RULE_W, Inches(0.65),
        runs=[("Firms that operate in the inelastic part of demand ",
               {'bold': True}),
              ("should raise their price", {'bold': True, 'color': RED})],
        size=24, fill_rgb=RED, fill_alpha_pct=12, border=RED)
    # 2026-08-26 (Nico): "elaborated in Videos 1+2" leaves the title and
    # becomes a navy video pointer beside the intuition bullet, above
    # the Uber shot
    _add_rounded_filled_box(
        slide, Inches(9.09), Inches(2.35), Inches(3.56), Inches(0.95),
        "%s  Elaborated in\nModule 2, Videos 1+2" % PV_GLYPH,
        fill=NAVY, text_color=WHITE, size=18, bold=True,
        corner_pct=0.15)
    # 2026-08-25 (Nico): his Uber shot, with the deck's rounded corners
    # and soft shade
    _add_media_image(slide, "uber_app.png", left=Inches(9.09),
                     top=Inches(5.18), width=Inches(3.56))
    return slide


# --------------------------------------------------------------------------
# Slide 54 — Estimated demand curve for Uber (research figure, kept as image)
# --------------------------------------------------------------------------

def slide_54_uber(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_OWN)
    _draw_action_title(
        slide, "Estimated Demand Curve for Uber (Surge Prices, 2016)")
    _add_text(slide, MARGIN, Inches(1.38), RULE_W, Inches(0.35),
              "I’ll provide detail on this estimation in my office hour – "
              "not relevant for the exam", size=15, italic=True,
              color=GRAY, font="Calibri")
    _add_media_image(slide, "image49.emf", left=Inches(2.05),
                     top=Inches(1.95), width=Inches(6.6),
                     rounded=False, shadow=True)
    # 2026-08-26 (Nico): wider, flatter callout ...
    _add_convention_box(
        slide, Inches(9.00), Inches(3.802), Inches(3.756), Inches(0.58),
        runs=[("Elasticity = −0.4", {'bold': True}),
              (" at 2016 base price", {})], size=16)
    # ... and its pointer lands where it did in his ORIGINAL deck (old
    # slide 51): the head sat at 79.8% across / 67.6% down the chart
    # image, which on our copy of the figure is (7.32, 5.23)
    _add_arrow(slide, (Inches(8.95), Inches(4.30)),
               (Inches(7.32), Inches(5.23)), color=GOLD, weight_pt=2.0,
               head=True)
    # 2026-08-26 (Nico): the Freakonomics screenshot gets the deck's
    # rounded corners + shade, and the podcast link from the original
    # deck is restored (click opens; the ScreenTip names the episode)
    pod = _add_media_image(slide, "image52.png", left=Inches(9.069),
                           top=Inches(5.568), width=Inches(3.756),
                           rounded=True, shadow=True)
    _link_shape_to_url(
        slide, pod,
        "https://freakonomics.com/podcast/why-uber-is-an-economists-dream/",
        tooltip="Freakonomics Radio 258: Why Uber Is an Economist’s Dream")
    _draw_footer(slide, FOOTER_TEXT, 50)
    return slide


# --------------------------------------------------------------------------
# Slide 55 — Special cases of price elasticities
# --------------------------------------------------------------------------

def slide_55_special_cases(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_OWN)
    _draw_action_title(slide, "Special Cases of Price Elasticities")
    # 2026-08-26 (Nico): the whole slide lifts 0.30" to clear the new
    # takeaway pill at the bottom; both panel labels move above their
    # figures and the axis letters sit closer to their axes
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.15), Inches(1.40), Inches(12.4),
        Inches(0.55),
        [("Elasticity changes along the demand curve, except ...", 0)],
        size=24)
    # left: perfectly elastic (horizontal D)
    figL = SimpleFig(1.95, 6.05, 4.0, 3.4, 10, 10)
    _fig_axes(slide, figL)
    _add_text(slide, Inches(figL.l - 0.30), Inches(2.15),
              Inches(0.7), Inches(0.35), "P", size=20, bold=True,
              italic=True, color=NAVY, font="Calibri")
    _add_text(slide, Inches(figL.l + figL.w + 0.05),
              Inches(figL.b + 0.05), Inches(0.7), Inches(0.35), "Q",
              size=20, bold=True, italic=True, color=NAVY, font="Calibri")
    _add_arrow(slide, (figL.x(0.4), figL.y(5.2)), (figL.x(9.4), figL.y(5.2)),
               color=RED, weight_pt=3.0, head=False)
    _add_text(slide, figL.x(9.0), figL.y(6.6), Inches(0.5), Inches(0.4),
              "D", size=22, bold=True, color=RED, font="Calibri")
    _add_rounded_filled_box(
        slide, Inches(2.44), Inches(2.145), Inches(3.3), Inches(0.55),
        "Perfectly Elastic", fill=NAVY, text_color=WHITE, size=18,
        bold=True, corner_pct=0.15)
    _add_math_equation(
        slide, Inches(2.74), Inches(2.82), Inches(2.6), Inches(0.6),
        _oED() + _omml_text(' = −∞'), size_pt=20, color=NAVY)
    # right: perfectly inelastic (vertical D)
    figR = SimpleFig(7.85, 6.05, 4.0, 3.4, 10, 10)
    _fig_axes(slide, figR)
    _add_text(slide, Inches(figR.l - 0.30), Inches(2.175),
              Inches(0.7), Inches(0.35), "P", size=20, bold=True,
              italic=True, color=NAVY, font="Calibri")
    _add_text(slide, Inches(figR.l + figR.w + 0.05),
              Inches(figR.b + 0.05), Inches(0.7), Inches(0.35), "Q",
              size=20, bold=True, italic=True, color=NAVY, font="Calibri")
    _add_arrow(slide, (figR.x(4.8), figR.y(0.3)), (figR.x(4.8), figR.y(9.4)),
               color=RED, weight_pt=3.0, head=False)
    _add_text(slide, figR.x(5.2), figR.y(9.2), Inches(0.5), Inches(0.4),
              "D", size=22, bold=True, color=RED, font="Calibri")
    _add_rounded_filled_box(
        slide, Inches(8.25), Inches(2.17), Inches(3.3), Inches(0.55),
        "Perfectly Inelastic", fill=NAVY, text_color=WHITE, size=18,
        bold=True, corner_pct=0.15)
    _add_math_equation(
        slide, Inches(7.90), Inches(2.732), Inches(2.4), Inches(0.6),
        _oED() + _omml_text(' = 0'), size_pt=20, color=NAVY)
    # 2026-08-26 (Nico): the gold takeaway pill he copied from CT's deck
    _add_rounded_filled_box(
        slide, Inches(3.58), Inches(6.596), Inches(6.85), Inches(0.399),
        "Elasticity depends on the starting price — except in these "
        "two extremes",
        fill=GOLD, text_color=NAVY, size=16, bold=True, italic=True,
        corner_pct=0.18)
    _draw_footer(slide, FOOTER_TEXT, 51)
    return slide


# --------------------------------------------------------------------------
# Slide 56 — What determines the price elasticity of demand?
# --------------------------------------------------------------------------

def slide_56_determinants(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_OWN)
    _draw_action_title(
        slide, "What Determines the Price Elasticity of Demand?")
    # 2026-08-26 (Nico): bigger type (28/24) and three empty paragraphs
    # that open a gap under "Closeness" for his four substitute photos
    _add_hierarchical_bullets(
        slide, Inches(0.383), Inches(1.494), Inches(5.3), Inches(5.6),
        [
            ("Availability of substitutes", 0),
            ("Number", 1),
            ("Closeness", 1),
            # 2026-08-26 (Nico, 2nd pass): the two remaining factors move
            # UP under the substitutes block and the photo gap drops to
            # the bottom of the column
            ("Willingness to switch", 0, {'space_before_pts': 6}),
            ("", 0),
            ("", 0),
            ("", 0),
            ([("Firm size relative to relevant market ", {}),
              ("(smaller size → higher demand elasticity)",
               {'size': 20})], 0, {'space_before_pts': 6}),
        ],
        size=28, sub_size=24)
    # his four jars: the same product in near-substitute variants, the
    # personalised "nico" jar last and largest (rounded + shaded, the
    # deck's photo treatment; the three small ones stay flat as pasted)
    for fname, x, y, w in (
            ("subs_nutella.png", 0.355, 3.560, 1.383),
            ("subs_skippy.png", 2.040, 3.560, 1.025),
            ("subs_nutella_peanut.png", 3.398, 3.560, 1.400)):
        _add_media_image(slide, fname, left=Inches(x), top=Inches(y),
                         width=Inches(w), rounded=False, shadow=False)
    _add_media_image(slide, "subs_nico_nutella.png", left=Inches(4.995),
                     top=Inches(1.970), width=Inches(2.804),
                     rounded=True, shadow=True)
    # 2026-08-26 (Nico): the elasticity-estimates table was still the
    # original screenshot (~8 pt type).  Rebuilt as a native table at
    # 16 pt, section rows merged across both columns — and (later the
    # same day) split into its THREE sections so they can be revealed
    # one at a time.  ONE backing card carries the shade for all three;
    # three stacked cards would cast shadows across each other.
    T_X, T_W, T_TOP = Inches(8.125), Inches(4.80), 1.45
    ROW_H, HDR_H = 0.311, 0.60
    COLS = [Inches(2.60), Inches(2.20)]
    _add_graphicframe_shadow(slide, T_X - Inches(0.15),
                             Inches(T_TOP - 0.15), T_W + Inches(0.30),
                             Inches(HDR_H + 15 * ROW_H + 0.30))
    sections = [
        # (rows, header row?)  the first block carries the column header
        ([["Product", "Estimated Own-Price Elasticity of Demand"],
          ["BROAD FOOD GROUPS", ""],
          ["Eggs", "−0.06"],
          ["Beef", "−0.35"],
          ["Fish", "−0.39"],
          ["Juice", "−1.05"]], True),
        ([["SPECIFIC BREAKFAST CEREALS", ""],
          ["Cap’N Crunch", "−2.28"],
          ["Froot Loops", "−2.34"],
          ["Kellogg’s Corn Flakes", "−3.38"],
          ["Cheerios", "−3.66"],
          ["Shredded Wheat", "−4.25"]], False),
        ([["SPECIFIC AUTOMOBILES", ""],
          ["Jeep Grand Cherokee", "−3.06"],
          ["Cadillac Seville", "−3.16"],
          ["Toyota Corolla", "−3.92"]], False),
    ]
    y = T_TOP
    for rows, has_header in sections:
        n_body = len(rows) - (1 if has_header else 0)
        heights = ([Inches(HDR_H)] if has_header else []) \
            + [Inches(ROW_H)] * n_body
        sect_row = 1 if has_header else 0     # the SECTION label row
        # the section label keeps the body look even in row 0, where
        # _add_styled_table would otherwise paint a navy header
        fills = {(sect_row, c): CREAM for c in (0, 1)}
        colors = {(sect_row, c): NAVY for c in (0, 1)}
        h = sum(heights) / 914400.0
        gf = _add_styled_table(
            slide, T_X, Inches(y), T_W, Inches(h), rows,
            col_widths=COLS, row_heights=heights,
            font_size=16, header_size=16, margin_v=Inches(0.02),
            first_col_bold=False, first_col_align_left=False,
            cell_fills=fills, cell_text_colors=colors, backing=False)
        tbl = gf.table
        cell = tbl.cell(sect_row, 0)          # label spans both columns
        cell.merge(tbl.cell(sect_row, 1))
        tf = cell.text_frame
        for extra in list(tf.paragraphs[1:]):
            extra._p.getparent().remove(extra._p)
        for run in tf.paragraphs[0].runs:
            run.font.bold = True
        y += h
    _add_text(slide, Inches(8.125), Inches(6.912), Inches(4.80),
              Inches(0.28),
              "Source: food data – Huang & Lin (2000), USDA Technical "
              "Bulletin 1887", size=11, italic=True, color=GRAY,
              font="Calibri")
    _draw_footer(slide, FOOTER_TEXT, 52)
    return slide


# --------------------------------------------------------------------------
# Slide 57 — Market vs. firm elasticity (NEW, from CT)
# --------------------------------------------------------------------------

def slide_57_market_vs_firm(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_OWN)
    _draw_action_title(slide, "Market vs. Firm Elasticity")
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.15), Inches(2.05), Inches(6.6),
        Inches(4.2),
        [
            ([("The elasticity a company faces is ", {}),
              ("not the same", {'bold': True}),
              (" as the elasticity the industry faces", {})], 0, {}),
            ("A firm's demand is more elastic when…", 0),
            ("the overall market elasticity is high", 1),
            ("the firm's market share is small", 1),
            ("competitors don't react to its price changes", 1),
        ],
        size=28, sub_size=24)
    # 2026-08-26 (Nico): his own aerial shot replaces CT's stock photo —
    # four competing stations on the four corners of one intersection,
    # which is exactly the market-vs-firm point
    # 2026-08-26 (Nico): the header counts the stations - "4 Gas
    # stations at ..."
    _add_text(slide, Inches(7.55), Inches(1.764), Inches(5.507),
              Inches(0.35),
              "4 Gas stations at Slauson Ave and La Brea Ave, "
              "Los Angeles", size=15,
              italic=True, bold=True, color=NAVY, font="Calibri",
              align=PP_ALIGN.CENTER)
    _add_media_image(slide, "Gas Stations Slauson.png",
                     left=Inches(7.55), top=Inches(2.09),
                     width=Inches(5.507))
    # 2026-08-26 (Nico): the point of the photo, right under it - folded
    # into the picture group by _group_pass's MANUAL_GROUPS_POST[53]
    _add_text(slide, Inches(7.55), Inches(6.30), Inches(5.507),
              Inches(0.50),
              "Demand at individual gas stations is much more elastic "
              "than the overall market demand for gas",
              size=15, italic=True, color=NAVY, font="Calibri",
              align=PP_ALIGN.CENTER)
    _set_notes(slide, (
        "The distinction that matters for pricing: the elasticity YOUR "
        "company faces is not the elasticity the industry faces. Demand "
        "for gasoline overall is quite inelastic — drivers still need "
        "to fill up when prices rise. But demand at one particular gas "
        "station is highly elastic: if it raises its price and the "
        "station across the intersection does not, customers just drive "
        "over there. A firm's demand is more elastic when the market "
        "elasticity is high, when its market share is small, and when "
        "competitors don't match its price changes. The photo is the "
        "intersection of Slauson and La Brea in Los Angeles, with four "
        "competing stations on the four corners — the posted prices are "
        "readable from the street, so switching costs almost nothing."))
    _draw_footer(slide, FOOTER_TEXT, 53)
    return slide


# --------------------------------------------------------------------------
# Slide 58 — Next: other elasticities (three-types re-anchor, own-price done)
# --------------------------------------------------------------------------

def slide_58_other_elasticities(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_ELAST)
    _draw_action_title(slide, "Next: Other Elasticities")
    _three_types_cards(slide, dim={0})
    _add_text(slide, MARGIN, Inches(5.45), RULE_W, Inches(0.45),
              "Own-price elasticity ✓ — next: how demand responds to "
              "income, and to other goods’ prices",
              size=20, italic=True, color=GRAY, font="Calibri",
              align=PP_ALIGN.CENTER)
    _draw_footer(slide, FOOTER_TEXT, 54)
    return slide




# ==========================================================================
#  BATCH E — income + cross-price elasticity, cheat sheet, wrap-up (59–76)
# ==========================================================================

def _add_polyline(slide, pts, *, color=NAVY, weight_pt=2.5, dash=None):
    """Open freeform polyline through pts [(x_emu, y_emu), ...]."""
    fb = slide.shapes.build_freeform(int(pts[0][0]), int(pts[0][1]),
                                     scale=1.0)
    fb.add_line_segments([(int(x), int(y)) for x, y in pts[1:]],
                         close=False)
    shp = fb.convert_to_shape()
    shp.fill.background()
    shp.shadow.inherit = False
    shp.line.color.rgb = color
    shp.line.width = Pt(weight_pt)
    if dash:
        ln = shp._element.spPr.find(qn('a:ln'))
        d = ln.makeelement(qn('a:prstDash'), {'val': dash})
        ln.append(d)
    return shp


# --------------------------------------------------------------------------
# Slide 59 — Income elasticity (definition)
# --------------------------------------------------------------------------

def slide_59_income_elasticity(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_INCOME)
    _draw_action_title(slide, "Income Elasticity")
    # 2026-08-26 (Nico): this slide is now CT's In-Class slide 47 — the
    # definition AND the categories on one page, in our formatting.  It
    # absorbs the old "Income Elasticity: Categories" slide, which is
    # deleted (his instruction the same day).
    _add_math_equation(
        slide, Inches(4.35), Inches(1.45), Inches(4.6), Inches(1.25),
        _oEI() + _omml_text(' = ')
        + _omml_frac(_o_pct('Q'), _omml_text('%Δ') + _omml_run('I')),
        size_pt=32, color=NAVY, fill=CREAM, line=NAVY, rounded=True,
        shadow=True)
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.15), Inches(2.80), Inches(12.4),
        Inches(2.0),
        [
            ([("% change in quantity demanded, divided by the % change "
               "in income (", {}), ("I", {'italic': True}),
              (")", {})], 0, {}),
            ([("Intuition:", {'bold': True}),
              ("  How sensitive demand is to changes in customer’s "
               "income", {})], 0, {}),
            ([("If  ", {}), ("E", {'italic': True}), ("I", {'size': 16}),
              (" > 0  →  ", {}),
              ("normal good", {'bold': True, 'color': CBLUE}),
              ("   (if  ", {}), ("E", {'italic': True}),
              ("I", {'size': 16}), (" > 1  →  ", {}),
              ("luxury good", {'bold': True, 'color': CBLUE}),
              (")", {})], 0, {}),
            ([("If  ", {}), ("E", {'italic': True}), ("I", {'size': 16}),
              (" < 0  →  ", {}),
              ("inferior good", {'bold': True, 'color': CBLUE})], 0, {}),
        ],
        size=24, line_spacing_pts=12)
    # her three cars, in our picture treatment (rounded + shade) with a
    # caption each, ordered to follow the two bullets above
    for fname, x, w, label in (
            ("ct_income_normal.png", 1.600, 2.773, "normal"),
            ("ct_income_luxury.png", 4.923, 3.264, "luxury"),
            ("ct_income_inferior.png", 8.737, 2.996, "inferior")):
        # 2026-08-26 (Nico): cars nudged down 0.07" (was 5.00 / 6.60)
        # and the labels enlarged to 16 pt bold (was 13 pt regular)
        _add_media_image(slide, fname, left=Inches(x), top=Inches(5.07),
                         width=Inches(w))
        _add_text(slide, Inches(x), Inches(6.67), Inches(w),
                  Inches(0.28), label, size=16, italic=True, bold=True,
                  color=GRAY, font="Calibri", align=PP_ALIGN.CENTER)
    _draw_footer(slide, FOOTER_TEXT, 55)
    return slide


# --------------------------------------------------------------------------
# Slide 60 — Example: income elasticity (Rivian R3)
# --------------------------------------------------------------------------

def slide_60_rivian(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_INCOME)
    _draw_action_title(slide, "Example: Income Elasticity")
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.15), Inches(1.85), Inches(12.4),
        Inches(1.9),
        [
            ("Average income in the U.S. goes up by 2%", 0),
            ("As a result, U.S. demand for the Rivian R3 goes up by 5%",
             0),
            ("What is the implied income elasticity of demand for the "
             "Rivian R3?", 0, {'bold': True}),
        ],
        size=26)
    _add_media_image(slide, "image54.png", left=Inches(3.95),
                     top=Inches(3.90), width=Inches(4.4))
    _draw_footer(slide, FOOTER_TEXT, 56)
    _add_pollbreak_badge(slide)
    return slide


# --------------------------------------------------------------------------
# Slide 63 — Solution: income elasticity of R3
# --------------------------------------------------------------------------

def slide_63_rivian_solution(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_INCOME)
    _draw_action_title(slide, "Solution: Income Elasticity of R3")
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.15), Inches(1.85), Inches(12.4),
        Inches(2.2),
        [
            ([("Answer is ", {}), ("2.5", {'bold': True})], 0, {}),
            ([("Recall that  ", {}), ("Eᵢ = %ΔQ / %ΔI", {'italic': True})],
             0, {}),
            ("We know:", 0),
            ("%ΔQ = 5%", 1),
            ("%ΔI = 2%", 1),
        ],
        size=24, sub_size=22, line_spacing_pts=12)
    _add_math_equation(
        slide, Inches(3.85), Inches(4.55), Inches(5.6), Inches(1.1),
        _oEI() + _omml_text(' = ')
        + _omml_frac(_omml_text('5%'), _omml_text('2%'))
        + _omml_text(' = 2.5'),
        size_pt=28, color=RED)
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.15), Inches(5.95), Inches(12.4),
        Inches(0.55),
        [("Does this seem like a strong reaction of demand?", 0,
          {'italic': True, 'bullet_style': 'none'})],
        size=24)
    _draw_footer(slide, FOOTER_TEXT, 59)
    return slide


# --------------------------------------------------------------------------
# Slide 65 — Which retailers do well in a recession? (native rebuild of the
# Stevenson/Wolfers Target-vs-Walmart figure, decision 2026-08-14)
# --------------------------------------------------------------------------

def slide_65_recession_retailers(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_INCOME)
    _draw_action_title(
        slide, "Which Retailers Do Well in a Recession? "
               "Clues from Income Elasticity")
    # 2026-08-26 (Nico): the three descriptive lines move into the
    # speaker notes and the figure takes the whole slide instead
    # native chart: stock prices Dec 2006 – Dec 2008 (approx. from figure)
    fig = SimpleFig(2.9, 6.20, 8.0, 4.1, 24, 70)
    _add_graphicframe_shadow(slide, Inches(1.35), Inches(1.95),
                             Inches(10.9), Inches(4.80))
    _fig_axes(slide, fig)
    _add_text(slide, Inches(fig.l - 1.35), Inches(1.58),
              Inches(1.9), Inches(0.32), "Stock price", size=15,
              bold=True, italic=True, color=NAVY, font="Calibri")
    for v in (25, 40, 55):
        _fig_ytick(slide, fig, v, "$%d" % v, size=13)
    for m, lbl in ((0, "Dec. 2006"), (12, "Dec. 2007"), (24, "Dec. 2008")):
        _fig_xtick(slide, fig, m, lbl, size=13)
    target = [(0, 57), (2, 60), (4, 63), (6, 59), (8, 55), (10, 52),
              (12, 51), (14, 53), (16, 50), (18, 46), (20, 40), (22, 34),
              (24, 31)]
    walmart = [(0, 46), (2, 47), (4, 48), (6, 44), (8, 43), (10, 45),
               (12, 47), (14, 50), (16, 55), (18, 57), (20, 53), (22, 54),
               (24, 55)]
    _add_polyline(slide, [(fig.x(m), fig.y(v)) for m, v in target],
                  color=RED, weight_pt=2.5)
    _add_polyline(slide, [(fig.x(m), fig.y(v)) for m, v in walmart],
                  color=GREEN, weight_pt=2.5)
    _add_arrow(slide, (fig.x(12), fig.y(0)), (fig.x(12), fig.y(68)),
               color=GRAY, weight_pt=1.25, head=False, dash="dash")
    # the recession label moves to the LEFT of its line: everything to
    # the right of December 2007 is hidden until the class has guessed
    _add_text(slide, fig.x(12) - Inches(2.70), fig.y(67), Inches(2.60),
              Inches(0.3), "U.S. economy enters recession", size=12,
              italic=True, color=GRAY, font="Calibri",
              align=PP_ALIGN.RIGHT)
    # the two verdicts sit INSIDE the plot now, so the cover card can
    # stop at the backing card's edge and leave its shade intact
    _add_text(slide, fig.x(17.3), fig.y(64), Inches(2.2),
              Inches(0.35), "Walmart’s stock rose", size=14, bold=True,
              color=GREEN, font="Calibri")
    _add_text(slide, fig.x(17.3), fig.y(26), Inches(2.2),
              Inches(0.35), "Target’s stock fell", size=14, bold=True,
              color=RED, font="Calibri")
    # 2026-08-26 (Nico): the white cover card — drawn LAST so it sits on
    # top of both series and their labels.  It fades out on the single
    # click of this slide (plan 60: "x:osp:6"), which is the reveal.
    _add_rect(slide, Inches(6.92), Inches(1.96), Inches(5.33),
              Inches(4.23), WHITE)
    _add_text(slide, Inches(1.35), Inches(6.85), Inches(10.9),
              Inches(0.28),
              "Source: Stevenson/Wolfers, Principles of Economics, 1e, "
              "© 2020 Worth Publishers (Figure 9); series approximate",
              size=11, italic=True, color=GRAY, font="Calibri",
              align=PP_ALIGN.CENTER)
    _draw_footer(slide, FOOTER_TEXT, 60)
    return slide


# --------------------------------------------------------------------------
# Slide 66 — Inferior goods in the news (NEW, from CT: ThredUp / WSJ)
# --------------------------------------------------------------------------

def slide_66_inferior_news(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_INCOME)
    _draw_action_title(slide, "Inferior Goods in the News")
    # 2026-08-26 (Nico): the WSJ masthead sits above its headline, the
    # photo is bigger, and its source line moved underneath it (so the
    # grouping pass pairs photo + source into one object)
    _add_media_image(slide, "ct_inferior_image28.png", left=Inches(0.95),
                     top=Inches(2.30), width=Inches(6.1),
                     rounded=False, shadow=False)
    _add_media_image(slide, "ct_inferior_image29.png", left=Inches(2.30),
                     top=Inches(1.873), width=Inches(2.6),
                     rounded=False, shadow=False)
    _add_media_image(slide, "ct_inferior_image27.jpg", left=Inches(7.46),
                     top=Inches(1.816), width=Inches(5.27))
    _add_text(slide, Inches(8.956), Inches(5.434), Inches(3.044),
              Inches(0.28),
              "Source: The Wall Street Journal (2025)", size=12,
              italic=True, color=GRAY, font="Calibri")
    # 2026-08-26 (Nico): the bottom-left corner was empty — the takeaway
    # this clipping is here to make goes in it
    _add_convention_box(
        slide, Inches(0.95), Inches(4.75), Inches(6.10), Inches(1.55),
        # 2026-08-28 (Nico's hand edit): he rewrote this line - it
        # leads with the recession framing instead of the
        # "inferior ≠ low quality" caveat, and "rises" carries the
        # emphasis (was: "Inferior ≠ low quality: demand RISES when
        # incomes fall. …")
        runs=[("Inferior goods and recessions: ", {'bold': True}),
              ("Demand ", {}),
              ("rises", {'bold': True}),
              (" when incomes fall. Secondhand clothing, "
               "discount grocers, store brands and public transit all "
               "grow in a downturn", {})],
        size=20)
    _set_notes(slide, (
        "A current example of an inferior good: ThredUp, the online "
        "secondhand-clothing marketplace. When consumers get worried "
        "about the economy and incomes tighten, demand for secondhand "
        "clothes rises — the number of active buyers rose 17% to 1.47 "
        "million in the second quarter. That is the signature of an "
        "inferior good: demand moves opposite to income."))
    _draw_footer(slide, FOOTER_TEXT, 61)
    return slide


# --------------------------------------------------------------------------
# Slide 67 — Cross-price elasticity (definition)
# --------------------------------------------------------------------------

def slide_67_cross_price(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_CROSS)
    _draw_action_title(slide, "Cross-Price Elasticity")
    _add_math_equation(
        slide, Inches(1.15), Inches(1.80), Inches(4.6), Inches(1.35),
        _oEXY() + _omml_text(' = ')
        + _omml_frac(_omml_text('%Δ')
                     + _omml_sub(_omml_run('Q'), _omml_run('X')),
                     _omml_text('%Δ')
                     + _omml_sub(_omml_run('P'), _omml_run('Y'))),
        size_pt=30, color=NAVY, fill=CREAM, line=NAVY, rounded=True,
        shadow=True)
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.15), Inches(3.55), Inches(8.4),
        Inches(3.2),
        [
            ("% change in the quantity demanded of good X, divided by "
             "the % change in the price of another good Y", 0),
            ([("Intuition:", {'bold': True}),
              ("  How sensitive demand is to changes in the price of "
               "another good", {})], 0, {}),
            ("Two cases:", 0),
            ([("If  ", {}), ("Eₓ,ᵧ > 0", {'italic': True}),
              (" ,  then the goods are ", {}),
              ("substitutes", {'bold': True,
                               'color': RGBColor(0x00, 0x70, 0xC0)})],
             1, {}),
            ([("If  ", {}), ("Eₓ,ᵧ < 0", {'italic': True}),
              (" ,  then the goods are ", {}),
              ("complements", {'bold': True,
                               'color': RGBColor(0x00, 0x70, 0xC0)})],
             1, {}),
        ],
        size=24, sub_size=22)
    # example pictures: substitutes (two ice-cream brands), complements
    _add_media_image(slide, "image63.png", left=Inches(9.55),
                     top=Inches(3.30), width=Inches(1.55))
    _add_media_image(slide, "image65.png", left=Inches(11.25),
                     top=Inches(3.30), width=Inches(1.55))
    # 2026-08-26 (Nico): both labels move up (4.95 -> 4.75,
    # 6.95 -> 6.81) and the gas pump shrinks (11.45/5.30/1.15" ->
    # 11.58/5.36/0.89")
    _add_text(slide, Inches(9.55), Inches(4.75), Inches(3.25),
              Inches(0.3), "substitutes", size=13, italic=True,
              color=GRAY, font="Calibri", align=PP_ALIGN.CENTER)
    _add_media_image(slide, "image63.png", left=Inches(9.55),
                     top=Inches(5.35), width=Inches(1.55))
    _add_media_image(slide, "image66.png", left=Inches(11.58),
                     top=Inches(5.36), width=Inches(0.89))
    _add_text(slide, Inches(9.55), Inches(6.81), Inches(3.25),
              Inches(0.3), "complements", size=13, italic=True,
              color=GRAY, font="Calibri", align=PP_ALIGN.CENTER)
    _draw_footer(slide, FOOTER_TEXT, 62)
    return slide


# --------------------------------------------------------------------------
# Slide 68 — Example: movie ticket and popcorn
# --------------------------------------------------------------------------

def slide_68_popcorn(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_CROSS)
    _draw_action_title(slide, "Example: Movie Ticket and Popcorn")
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.15), Inches(1.85), Inches(12.4),
        Inches(1.6),
        [
            ([("When a movie theater ", {}),
              ("increased", {'italic': True}),
              (" its ticket price from $15 to $18, the quantity demanded "
               "of popcorn ", {}), ("decreased", {'italic': True}),
              (" by 8%", {})], 0, {}),
            ("What is the implied cross-price elasticity?", 0,
             {'bold': True}),
        ],
        size=26)
    _add_media_image(slide, "image67.jpeg", left=Inches(3.85),
                     top=Inches(3.55), width=Inches(5.6))
    _draw_footer(slide, FOOTER_TEXT, 63)
    _add_pollbreak_badge(slide)
    return slide


# --------------------------------------------------------------------------
# Slide 71 — Solution (popcorn)
# --------------------------------------------------------------------------

def slide_71_popcorn_solution(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_CROSS)
    _draw_action_title(slide, "Solution")
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.15), Inches(1.85), Inches(12.4),
        Inches(2.0),
        [
            ([("Unit free  →  ", {}), ("not", {'bold': True}),
              (" 5%", {})], 0, {}),
            ("Change in quantity of good X:   %ΔQₓ = −8%", 0),
            ("Change in price of good Y:   %ΔPᵧ = (18 − 15)/15 = +20%",
             0),
        ],
        size=24, line_spacing_pts=14)
    _add_math_equation(
        slide, Inches(3.55), Inches(4.15), Inches(6.2), Inches(1.15),
        _oEXY() + _omml_text(' = ')
        + _omml_frac(_omml_text('−8%'), _omml_text('+20%'))
        + _omml_text(' = −0.4'),
        size_pt=28, color=RED)
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.15), Inches(5.65), Inches(12.4),
        Inches(1.0),
        [([("Popcorn and movie tickets are …  ", {}),
           ("…complements", {'bold': True,
                             'color': RGBColor(0x00, 0x70, 0xC0)})], 0,
          {})],
        size=26)
    _draw_footer(slide, FOOTER_TEXT, 66)
    return slide


# --------------------------------------------------------------------------
# Slide 72 — Cross-price elasticity in the news (NEW, from CT: hybrids)
# --------------------------------------------------------------------------

def slide_72_crossprice_news(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_CROSS)
    _draw_action_title(slide, "Cross-Price Elasticity in the News")
    # 2026-08-26 (Nico): his layout - the masthead moves ABOVE the
    # headline (was 1.75 / 5.35), the clipping grows to 6.83" (was
    # 6.0" at x 0.95) and the photo to 5.26" (was 4.9" at 7.65/2.05).
    # He also deleted the "Source: The Wall Street Journal" line; the
    # masthead itself carries the attribution now.
    _add_media_image(slide, "ct_crossprice_image34.png",
                     left=Inches(0.61), top=Inches(2.35),
                     width=Inches(6.83), rounded=False, shadow=False)
    _add_media_image(slide, "ct_crossprice_image29.png",
                     left=Inches(2.48), top=Inches(1.75),
                     width=Inches(2.6), rounded=False, shadow=False)
    # 2026-08-28 (Nico): the right-hand column now carries BOTH sides of
    # the cross-price story, stacked - the fuel-hungry truck on top (the
    # complement whose demand falls) and the EV that demand moves to
    # underneath.  This replaces the single Toyota plug-in photo
    # (ct_crossprice_image35.png at 7.56 / 1.97, 5.26" wide), which the
    # slide had no text to explain.  Both photos are 16:9 so they stack
    # cleanly; the F-250 shot was cropped from 4:3 to 16:9 on download.
    # 2026-08-28 (Nico's hand edit): he replaced BOTH photos with two of
    # his own and deleted the "Photos: Wikimedia Commons" caption.  His
    # pictures and positions are adopted verbatim (they sit at their
    # native aspect ratios); the deck's photo treatment - rounded
    # corners + soft drop shadow - is applied on top, since he inserted
    # them raw.  (An earlier round used two Wikimedia photos here; those
    # files were removed at cleanup on 2026-08-28.)
    # 2026-08-28 (Nico, second pass): the top photo is a Gulf station
    # with a pickup and an Escalade at the pumps, replacing the
    # close-up nozzle shot (which carried 123RF watermarks)
    _add_media_image(slide, "nv_gulf_station.png",
                     left=Inches(7.995), top=Inches(1.573),
                     width=Inches(4.963), height=Inches(2.532))
    _add_media_image(slide, "nv_rivian_charger.png",
                     left=Inches(8.294), top=Inches(4.402),
                     width=Inches(4.531), height=Inches(2.690))
    # 2026-08-26 (Nico): the slide's message, in the empty lower-left
    # 2026-08-28 (Nico): a second line explains the EV photo - demand
    # moves to the substitute once electricity is relatively cheaper.
    # Split into TWO boxes the same evening, so the complements point
    # and the substitution point can be revealed on separate clicks.
    # Each pairs with the photo above it in the build: the Gulf station
    # with the complements box, the Rivian with the substitute box.
    _add_convention_box(
        slide, Inches(0.61), Inches(4.12), Inches(6.83), Inches(1.20),
        runs=[
            ("As the price of fuel goes up, demand for fuel-powered "
             "vehicles goes down", {}),
            ("\u2192  The two are ", {'newline': True}),
            ("complements", {'bold': True,
                             'color': RGBColor(0x00, 0x70, 0xC0)}),
        ],
        size=20, align=PP_ALIGN.CENTER, line_spacing_pct=110)
    _add_convention_box(
        slide, Inches(0.61), Inches(5.50), Inches(6.83), Inches(1.20),
        runs=[
            ("Demand then shifts to a ", {}),
            ("substitute", {'bold': True,
                            'color': RGBColor(0x00, 0x70, 0xC0)}),
            (" for fuel-powered vehicles: electric vehicles, as "
             "electricity has become ", {}),
            ("relatively", {'bold': True}),
            (" cheaper", {}),
        ],
        size=20, align=PP_ALIGN.CENTER, line_spacing_pct=110)
    _set_notes(slide, (
        "Cross-price elasticity in action: when gasoline prices spike, "
        "sales of hybrid cars surge. Gasoline and conventional cars are "
        "complements — expensive gas makes gas-guzzlers less attractive "
        "— while hybrids are a substitute for conventional cars. The "
        "sticker shock at the pump shows up directly in the showroom. "
        "The two photos make the same point: the fuel-hungry pickup on "
        "top is the complement whose demand falls, and the electric "
        "truck below is the substitute demand moves to.\n\n"
        "Photo credits (Wikimedia Commons): top — "
        "File:19 Ford F-250 Super Duty XLT.jpg, HJUdall, CC0, cropped "
        "to 16:9; bottom — File:EV Charging Station (53857454477).jpg, "
        "ajay_suresh, CC BY 2.0 (a Rivian R1T at a charging station)."))
    _draw_footer(slide, FOOTER_TEXT, 67)
    return slide


# --------------------------------------------------------------------------
# Slide 73 — Estimated cross-price elasticities: ready-to-eat cereal
# (Nevo 2001; native table + original cereal-box images)
# --------------------------------------------------------------------------

def slide_73_cereal(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_CROSS)
    _draw_action_title(
        slide, "Estimated Cross-Price Elasticities: "
               "Ready-to-Eat Cereal")
    rows = [
        ["", "Corn Flakes", "Frosted Flakes", "Rice Krispies"],
        ["Frosted Flakes", "0.15", "", ""],
        ["Rice Krispies", "0.19", "0.14", ""],
        ["Froot Loops", "0.02", "0.13", "0.04"],
    ]
    # 2026-08-26 (Nico): the table's card gets rounded edges (the
    # cereal-box images stay flat)
    _add_styled_table(slide, Inches(3.05), Inches(2.45), Inches(7.4),
                      Inches(3.1), rows, font_size=18, header_size=18,
                      backing_rounded=True)
    # cereal boxes: columns (above header) and rows (left of first column)
    for fname, x, y in (("image70.png", Inches(5.47), Inches(1.62)),
                        ("image72.png", Inches(7.32), Inches(1.62)),
                        ("image73.png", Inches(9.17), Inches(1.62))):
        _add_media_image(slide, fname, left=x, top=y, width=Inches(0.72),
                         rounded=False, shadow=False)
    for fname, y in (("image72.png", Inches(3.35)),
                     ("image73.png", Inches(4.15)),
                     ("image74.png", Inches(4.95))):
        _add_media_image(slide, fname, left=Inches(2.25), top=y,
                         width=Inches(0.68), rounded=False, shadow=False)
    _add_hierarchical_bullets(
        slide, Inches(3.05), Inches(5.95), Inches(7.4), Inches(0.5),
        [([("Larger cross-price elasticity  →  ", {}),
           ("closer substitutes", {'bold': True})], 0,
          {'bullet_style': 'none',
           'color': RGBColor(0x00, 0x70, 0xC0)})],
        size=20)
    _add_text(slide, Inches(1.35), Inches(6.78), Inches(10.6),
              Inches(0.3),
              "Source: Nevo, A. (2001). Measuring Market Power in the "
              "Ready-to-Eat Cereal Industry. Econometrica, 69(2), "
              "307–342", size=11, italic=True, color=GRAY,
              font="Calibri", align=PP_ALIGN.CENTER)
    _set_notes(slide, (
        "This table summarizes the cross-price elasticities between four "
        "products in this space: Frosted Flakes, Rice Krispies, Froot "
        "Loops, and Corn Flakes. Take the top left value for example. It "
        "means that the cross-price elasticity between Frosted Flakes "
        "and Corn Flakes is 0.15. In other words, if the price of Corn "
        "Flakes were to increase by 10%, the quantity demanded of "
        "Frosted Flakes would rise by 1.5%. This suggests that Frosted "
        "Flakes and Corn Flakes are substitutes."))
    _draw_footer(slide, FOOTER_TEXT, 68)
    return slide


# --------------------------------------------------------------------------
# Slide 74 — Cheat sheet: own-price elasticity (NEW, from CT)
# --------------------------------------------------------------------------

def slide_74_cheatsheet(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_WRAP)
    _draw_action_title(slide, "Cheat Sheet: Own-Price Elasticity")
    # left column: definition + categories
    _add_rounded_filled_box(slide, Inches(0.55), Inches(1.70),
                            Inches(5.9), Inches(0.52), "Definition",
                            fill=NAVY, text_color=WHITE, size=18,
                            bold=True, corner_pct=0.15)
    _add_math_equation(
        slide, Inches(0.55), Inches(2.40), Inches(5.9), Inches(1.25),
        _oED_frac('P'), size_pt=28, color=NAVY, fill=CREAM, line=NAVY,
        rounded=True, shadow=True)
    _add_hierarchical_bullets(
        slide, Inches(0.75), Inches(3.80), Inches(5.6), Inches(0.6),
        [([("Intuitively:", {'bold': True}),
           (" measures consumers’ price sensitivity", {})], 0,
          {'bullet_style': 'none'})],
        size=17)
    _add_rounded_filled_box(slide, Inches(0.55), Inches(4.55),
                            Inches(5.9), Inches(0.52), "Categories",
                            fill=NAVY, text_color=WHITE, size=18,
                            bold=True, corner_pct=0.15)
    _add_hierarchical_bullets(
        slide, Inches(0.75), Inches(5.25), Inches(5.6), Inches(1.7),
        [
            ([("Inelastic:", {'bold': True}),
              ("   −1 < Eᴅ < 0    (|Eᴅ| < 1)", {})], 0,
             {'bullet_style': 'none'}),
            ([("Unit-elastic:", {'bold': True}),
              ("   Eᴅ = −1", {})], 0, {'bullet_style': 'none'}),
            ([("Elastic:", {'bold': True}),
              ("   Eᴅ < −1    (|Eᴅ| > 1)", {})], 0,
             {'bullet_style': 'none'}),
        ],
        size=17, line_spacing_pts=8)
    # right column: computing
    _add_rounded_filled_box(slide, Inches(6.85), Inches(1.70),
                            Inches(5.9), Inches(0.52),
                            "Computing elasticity", fill=NAVY,
                            text_color=WHITE, size=18, bold=True,
                            corner_pct=0.15)
    _add_hierarchical_bullets(
        slide, Inches(7.05), Inches(2.35), Inches(5.6), Inches(0.4),
        [("Two observed points (approximation):", 0,
          {'bullet_style': 'none', 'bold': True})],
        size=16)
    _add_math_equation(
        slide, Inches(7.05), Inches(2.75), Inches(5.5), Inches(1.45),
        _omml_sub(_omml_run('E'), _omml_run('d')) + _omml_text(' = ')
        + _omml_frac(
            _omml_text('(') + _omml_sub(_omml_run('Q'), _omml_text('1'))
            + _omml_text('−')
            + _omml_sub(_omml_run('Q'), _omml_text('0'))
            + _omml_text(')/')
            + _omml_sub(_omml_run('Q'), _omml_text('0')),
            _omml_text('(') + _omml_sub(_omml_run('P'), _omml_text('1'))
            + _omml_text('−')
            + _omml_sub(_omml_run('P'), _omml_text('0'))
            + _omml_text(')/')
            + _omml_sub(_omml_run('P'), _omml_text('0'))),
        size_pt=22, color=NAVY, fill=CREAM, line=NAVY, rounded=True)
    _add_hierarchical_bullets(
        slide, Inches(7.05), Inches(4.40), Inches(5.6), Inches(0.4),
        [("Whole demand function known:", 0,
          {'bullet_style': 'none', 'bold': True})],
        size=16)
    _add_math_equation(
        slide, Inches(7.05), Inches(4.80), Inches(5.5), Inches(1.1),
        _omml_sub(_omml_run('E'), _omml_run('d')) + _omml_text(' = ')
        + _oD_slope() + _omml_text(' ∙ ')
        + _omml_frac(_omml_run('P'), _omml_run('Q')),
        size_pt=22, color=NAVY, fill=CREAM, line=NAVY, rounded=True)
    _add_hierarchical_bullets(
        slide, Inches(7.05), Inches(6.00), Inches(5.6), Inches(0.7),
        [("ΔQ/ΔP = slope of the demand function (Q as a function of P)",
          0, {'bullet_style': 'none', 'italic': True, 'color': GRAY})],
        size=14)
    _draw_footer(slide, FOOTER_TEXT, 69)
    return slide


# --------------------------------------------------------------------------
# Slides 75/76 — wrap-up outlines with post-work pointers
# --------------------------------------------------------------------------

def slide_75_postwork_videos(prs):
    # 2026-08-28 (Nico's hand edit): the title names the videos this
    # slide is pointing at (was the bare "Outline of Module 2").  The
    # tag is TAG_OUTLINE now - these are agenda slides.
    slide = make_m2_outline(prs, 70, section_tag=TAG_OUTLINE,
                            title="Outline of Module 2: Upcoming "
                                  "Videos 1+2",
                            highlight_set={2, 3, 4})
    # bottom-right link box overlaying the footer (deck convention),
    # drawn last so it sits in front
    _add_outlined_box(slide, Inches(8.15), Inches(6.68), Inches(4.9),
                      Inches(0.72),
                      "\u25b6  Module 2 Videos 1+2   \u00b7   "
                      "Practice Videos 1+2\nOn BL under "
                      "\u201cModule 2 Post-Work\u201d",
                      line=GOLD, text_color=NAVY, size=15, bold=True,
                      rounded=True, shadow=True, corner_pct=0.20)
    return slide


def slide_76_postwork_ps2(prs):
    # 2026-08-28 (Nico's hand edit): title names the upcoming video
    slide = make_m2_outline(prs, 71, section_tag=TAG_OUTLINE,
                            # 2026-08-28: he corrected this to Video 3,
                            # which is what the slide points at
                            title="Outline of Module 2: Upcoming "
                                  "Video 3",
                            highlight_set={5})
    # 2026-08-28 (Nico): the note moves off the right edge towards the
    # middle of the slide and becomes a WARNING box - the deck's
    # transparent dark-red tint with a dark-red border (same device as
    # slides 48 and 53) - so it stands out.  Its old corner position
    # (9.05 / 1.75, 3.9 x 1.5) also swallowed item 2's coverage pill,
    # which the grouping pass then paired with the box instead of the
    # note text; at 5.70 the box clears the pill column entirely.
    # Text goes 14 -> 18 pt (the deck's box-text floor).
    _add_convention_box(
        # 2026-08-28 (Nico's hand edit): pulled down and right a touch,
        # from 5.45 / 3.00, so it sits clear of item 3's row
        slide, Inches(5.55), Inches(3.55), Inches(5.50), Inches(1.50),
        runs=[("Note: ", {'bold': True, 'color': RED}),
              ("You do not need to perform the actual estimation "
               "(regression). But you need to understand how to "
               "interpret regression coefficients", {})],
        size=18, fill_rgb=RED, fill_alpha_pct=12, border=RED)
    _add_outlined_box(slide, Inches(8.15), Inches(5.30), Inches(4.9),
                      Inches(0.72),
                      "\u25b6  Module 2 Video 3   \u00b7   "
                      "\u270e  Problem Set 2\nOn BL under "
                      "\u201cModule 2 Post-Work\u201d",
                      line=GOLD, text_color=NAVY, size=15, bold=True,
                      rounded=True, shadow=True, corner_pct=0.20)
    return slide


# --------------------------------------------------------------------------
# build() — the 76-slide registry (approved outline, 2026-08-14).
# Slides start as make_todo placeholders and are swapped to real
# builders batch by batch; spliced slides stay make_stub until
# _splice_media.py replaces them.
# --------------------------------------------------------------------------

# --------------------------------------------------------------------------
# Deck-wide subscript pass (2026-08-24, Nico: "when there are formulas with
# a letter such as ED, the d is the subscript ... go through the whole deck
# and check that this is implemented").
#
# Two defects are fixed:
#   1. Unicode LOOKALIKES typed inline - the small-capital D in "E<D>", the
#      subscript digits in "P<0>" / "Q<1>", and friends.  They are not
#      subscripts, they are separate characters that happen to sit low.
#   2. A subscript FAKED with a smaller font: an "E" run followed by a "D"
#      run at a reduced size (this is what slide 9 was doing).
# Both are rewritten to a real PowerPoint subscript, i.e. the index run
# carries baseline="-25000", which is exactly what PowerPoint itself writes
# and what CT's own decks use.
# --------------------------------------------------------------------------

SUBSCRIPT_BASELINE = "-25000"

# lookalike -> the character it really stands for
SUB_LOOKALIKE = {
    "ᴅ": "D",     # small capital D, used for E_D
    "ᴄ": "C",     # small capital C
    "ᵢ": "i",     # subscript i, used for E_I
    "ᵧ": "y",     # subscript gamma, used for P_Y
    "ₓ": "x",     # subscript x, used for E_X / Q_X
}
for _d in range(10):
    SUB_LOOKALIKE[chr(0x2080 + _d)] = str(_d)

# a base letter immediately followed by one or more lookalikes
_SUB_RE = re.compile(
    "([A-Za-z])([" + "".join(SUB_LOOKALIKE) + "]+(?:,[" +
    "".join(SUB_LOOKALIKE) + "]+)*)")


def _sub_expand(idx_text):
    return "".join(SUB_LOOKALIKE.get(ch, ch) for ch in idx_text)


def _split_subscript_runs(para):
    """Split every run of a paragraph so an index becomes a real
    subscript run.  Formatting is copied from the base run; the italic
    state is left exactly as it was."""
    changed = False
    for run in list(para.runs):
        text = run.text or ""
        if not _SUB_RE.search(text):
            continue
        pieces = []                     # (text, is_subscript)
        pos = 0
        for m in _SUB_RE.finditer(text):
            if m.start() > pos:
                pieces.append((text[pos:m.start()], False))
            pieces.append((m.group(1), False))
            pieces.append((_sub_expand(m.group(2)), True))
            pos = m.end()
        if pos < len(text):
            pieces.append((text[pos:], False))
        r_el = run._r
        anchor = r_el
        for i, (txt, is_sub) in enumerate(pieces):
            if i == 0:
                run.text = txt
                continue
            new_r = copy.deepcopy(r_el)
            for t_el in new_r.findall(qn('a:t')):
                t_el.text = txt
            npr = new_r.find(qn('a:rPr'))
            if npr is not None:
                if is_sub:
                    npr.set('baseline', SUBSCRIPT_BASELINE)
                else:
                    npr.attrib.pop('baseline', None)
            anchor.addnext(new_r)
            anchor = new_r
        changed = True
    return changed


def _fix_smallfont_subscripts(para):
    """An index faked with a smaller font: a one-or-two-character run
    (D / I / X / a digit) that directly follows a run ending in a letter
    and is set smaller than it.  Give it the real baseline and put its
    size back to the base run's."""
    changed = False
    runs = list(para.runs)
    for i in range(1, len(runs)):
        prev, cur = runs[i - 1], runs[i]
        ptxt, ctxt = (prev.text or ""), (cur.text or "")
        if not ptxt or not ctxt:
            continue
        if not (ptxt[-1].isalpha() and len(ctxt) <= 2
                and (ctxt.isalpha() or ctxt.isdigit())):
            continue
        ppr = prev._r.find(qn('a:rPr'))
        cpr = cur._r.find(qn('a:rPr'))
        if ppr is None or cpr is None:
            continue
        if cpr.get('baseline'):
            continue                       # already a real subscript
        psz, csz = ppr.get('sz'), cpr.get('sz')
        if not psz or not csz or int(csz) >= int(psz):
            continue
        cpr.set('baseline', SUBSCRIPT_BASELINE)
        cpr.set('sz', psz)
        changed = True
    return changed


def _iter_text_frames(shp):
    if shp.has_text_frame:
        yield shp.text_frame
    if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shp.shapes:
            for tf in _iter_text_frames(child):
                yield tf
    if getattr(shp, "has_table", False) and shp.has_table:
        for row in shp.table.rows:
            for cell in row.cells:
                yield cell.text_frame


def apply_subscripts(prs):
    """Deck-wide: turn every faked index into a real subscript run."""
    n = 0
    for slide in prs.slides:
        for shp in slide.shapes:
            for tf in _iter_text_frames(shp):
                for para in tf.paragraphs:
                    if _split_subscript_runs(para):
                        n += 1
                    if _fix_smallfont_subscripts(para):
                        n += 1
    return n



def _add_external_link(slide, run, url, *, underline=True,
                       lock_color=True):
    """Make ``run`` open ``url`` in the browser.

    Same colour-lock trick as the slide-jump helper: without it
    PowerPoint repaints hyperlink text in the theme's hlink blue and the
    caption stops looking like a caption.
    """
    rId = slide.part.relate_to(url, RT.HYPERLINK, is_external=True)
    rPr = run._r.get_or_add_rPr()
    if underline:
        rPr.set('u', 'sng')
    for hl in rPr.findall(qn('a:hlinkClick')):
        rPr.remove(hl)
    hlinkClick = ET.SubElement(rPr, qn('a:hlinkClick'))
    hlinkClick.set(qn('r:id'), rId)
    if lock_color:
        AHYP_NS = ('http://schemas.microsoft.com/office/drawing/2018/'
                   'hyperlinkcolor')
        ext_xml = (
            f'<a:extLst xmlns:a="{A_NS}">'
            f'<a:ext xmlns:ahyp="{AHYP_NS}" '
            f'uri="{{A12FA001-AC4F-418D-AE19-62706E023703}}">'
            f'<ahyp:hlinkClr val="tx"/></a:ext></a:extLst>')
        hlinkClick.append(ET.fromstring(ext_xml))
    return run


# --------------------------------------------------------------------------
# Source links carried over from CT's decks (2026-08-24).  Keyed by the
# exact run text they belong on, so the pass finds them wherever the run
# ends up as slides are inserted or deleted.
# --------------------------------------------------------------------------

CT_SOURCE_LINKS = {
    "Source: Pharmaceutical Technology":
        "https://www.pharmaceutical-technology.com/news/trump-hits-long-"
        "term-pledge-as-us-prices-of-weight-loss-drugs-slashed/",
    "Novo Nordisk shares tumbled ~18%":
        "https://www.pharmaceutical-technology.com/news/novo-nordisk-"
        "shares-tumble-18-after-2026-sales-dip-warning/",
    "camelcamelcamel.com":
        "https://camelcamelcamel.com/product/B0DRVPMBX3",
    "tylervigen.com/spurious-correlations":
        "https://www.tylervigen.com/spurious-correlations",
}


def apply_ct_source_links(prs, links=None):
    """Deck-wide: attach the CT source URL to every run whose text matches
    one of CT_SOURCE_LINKS.  Runs that already carry a hyperlink are left
    alone."""
    links = CT_SOURCE_LINKS if links is None else links
    n = 0
    for slide in prs.slides:
        for shp in slide.shapes:
            for tf in _iter_text_frames(shp):
                for para in tf.paragraphs:
                    for run in para.runs:
                        url = links.get((run.text or "").strip())
                        if not url:
                            continue
                        rPr = run._r.find(qn('a:rPr'))
                        if rPr is not None and \
                                rPr.find(qn('a:hlinkClick')) is not None:
                            continue
                        _add_external_link(slide, run, url)
                        n += 1
    return n


# --------------------------------------------------------------------------
# Post-work reference boxes (2026-08-24, Nico).  Any pointer to a problem
# set or a teaching note is built here so the whole deck family uses one
# look and one glyph vocabulary:
#     PS_GLYPH   ✎   an exercise to work
#     TN_GLYPH   ▤   a note to read
# The label carries the problem-set NUMBER only - never the exercise
# numbers - so the reference survives re-numbering in later years.
# --------------------------------------------------------------------------

PS_GLYPH = "✎"
TN_GLYPH = "▤"
# 2026-08-25 (Nico): a pointer to a practice VIDEO gets the play
# glyph the deck already uses on its video link boxes
PV_GLYPH = "▶"
# 2026-08-26 (Nico): every problem-set pointer sits in the bottom-RIGHT
# corner, overlapping the footer — the slide-42 position.  Use this
# constant at every call site so the corner never drifts.
PS_BOX_XY = (Inches(10.17), Inches(6.53))


def _add_reference_box(slide, left, top, width, height, label, *,
                       kind="ps", size=15, corner_pct=0.25,
                       sub_label=None, sub_size=None):
    """Gold-bordered rounded pointer to post-work material.

    ``kind`` is "ps" (problem set) or "tn" (teaching note); it only picks
    the leading glyph.  Pass ``kind=None`` for a reference with no glyph.
    """
    glyph = {"ps": PS_GLYPH, "tn": TN_GLYPH, "video": PV_GLYPH}.get(kind)
    text = ("%s  %s" % (glyph, label)) if glyph else label
    return _add_outlined_box(slide, left, top, width, height, text,
                             line=GOLD, text_color=NAVY, size=size,
                             bold=True, rounded=True, shadow=True,
                             corner_pct=corner_pct,
                             sub_label=sub_label, sub_size=sub_size)


def build(out_path=None):
    prs = Presentation()
    prs.slide_width = int(SLIDE_W)
    prs.slide_height = int(SLIDE_H)
    # strip python-pptx default layouts down to one master (template rule)
    slide_01_title(prs)                                                #  1
    slide_02_logistics(prs)                      #  2
    slide_03_recap(prs)                  #  3
    make_stub(prs, 4, TAG_LOGISTICS, "Poll: pace feedback", STUB_POLL) #  4
    make_stub(prs, 5, TAG_LOGISTICS, "Poll results", STUB_POLL)        #  5
    slide_06_roadmap(prs)             #  6
    make_m2_outline(prs, 7, descriptions=True)                         #  7
    make_m2_outline(prs, 8, highlight_idx=0)                           #  8
    slide_09_law_of_demand(prs)                    #  9
    slide_10_more_customers(prs)   # 10
    make_stub(prs, 11, TAG_LAW, "Poll: WTP for pizza", STUB_POLL)      # 11
    make_stub(prs, 12, TAG_LAW, "Poll results", STUB_POLL)             # 12
    make_stub(prs, 13, TAG_LAW, "Class demand curve (live Excel)",
              STUB_EXCEL)                                              # 13
    slide_14_existing_buy_more(prs)      # 14
    slide_15_multiunit(prs)  # 15
    slide_16_gates(prs) # 16
    slide_17_inglehart(prs)        # 17
    slide_18_consumer_opt(prs)               # 18
    slide_19_movies(prs)            # 19
    slide_20_aggregation(prs)       # 20
    slide_21_factors(prs)            # 21
    slide_22_snob_news(prs)         # 22
    slide_23_network_effects(prs)  # 23
    slide_24_remember(prs)               # 24
    make_m2_outline(prs, 25, highlight_idx=1)                          # 25
    slide_26_generic_elasticity(prs)  # 26
    slide_27_netflix(prs)           # 27
    slide_28_what_is_elasticity(prs) # 28
    slide_29_three_types(prs)         # 29
    slide_30_own_price(prs)                # 30
    slide_31_water(prs)  # 31
    make_stub(prs, 32, TAG_OWN, "Poll results", STUB_POLL)             # 32
    slide_34_water_solution(prs)                     # 33
    slide_35_categories(prs)  # 34
    slide_36_yoga(prs)               # 35
    make_stub(prs, 36, TAG_OWN, "Poll results", STUB_POLL)             # 36
    slide_39_yoga_solution(prs)                      # 37
    slide_40_method1(prs)      # 38
    slide_45_megamillions(prs)             # 39
    make_stub(prs, 42, TAG_OWN, "Poll: Mega Millions elasticity",
              STUB_POLL)                   # 40
    slide_41_megamillions_solution(prs)    # 43
    slide_46_method2(prs)      # 42
    slide_47_point_steps(prs)      # 43
    slide_48_from_demand_fn(prs) # 44
    make_stub(prs, 45, TAG_OWN, "Poll: elasticity at P=2", STUB_POLL)  # 43
    make_stub(prs, 46, TAG_OWN, "Poll results", STUB_POLL)             # 44
    slide_51_qp_solution(prs)       # 47
    slide_52_linear_elasticity(prs)   # 48
    slide_53_insight(prs)                   # 49
    slide_54_uber(prs)     # 50
    slide_55_special_cases(prs) # 51
    slide_56_determinants(prs)  # 52
    slide_57_market_vs_firm(prs)          # 53
    slide_58_other_elasticities(prs)          # 54
    slide_59_income_elasticity(prs)                # 55
    slide_60_rivian(prs)  # 56
    make_stub(prs, 57, TAG_INCOME, "Poll: R3 income elasticity", STUB_POLL)   # 55
    make_stub(prs, 58, TAG_INCOME, "Poll results", STUB_POLL)          # 56
    slide_63_rivian_solution(prs)  # 59
    # slide 60 (Income Elasticity: Categories) deleted 2026-08-26 (Nico):
    # CT's slide 47, now our slide 55, carries the categories
    slide_65_recession_retailers(prs)  # 61
    slide_66_inferior_news(prs)       # 62
    slide_67_cross_price(prs)            # 63
    slide_68_popcorn(prs) # 64
    make_stub(prs, 64, TAG_CROSS, "Poll: popcorn cross-price", STUB_POLL)  # 63
    make_stub(prs, 65, TAG_CROSS, "Poll results", STUB_POLL)           # 64
    slide_71_popcorn_solution(prs)                 # 67
    slide_72_crossprice_news(prs)  # 68
    slide_73_cereal(prs)  # 69
    slide_74_cheatsheet(prs)  # 70
    slide_75_postwork_videos(prs)    # 71
    slide_76_postwork_ps2(prs)  # 72

    # CT's own source links, restored on the runs we adopted
    apply_ct_source_links(prs)

    # deck-wide subscript pass (2026-08-24): every faked index
    # becomes a real PowerPoint subscript run
    apply_subscripts(prs)

    # deck-wide speaker notes (2026-08-23): fill in every slide that does
    # not already carry notes of its own, so the substantive notes ported
    # verbatim from Nico's original deck are never overwritten.  Poll
    # slides are absent from NOTES on purpose — their notes ARE the
    # PollEverywhere mechanism — and slide 13's notes are injected by
    # _splice_media.py, since the splice replaces that slide wholesale.
    from _notes_m2 import NOTES as _M2_NOTES
    for _i, _slide in enumerate(prs.slides, start=1):
        _txt = _M2_NOTES.get(_i)
        if not _txt:
            continue
        if (_slide.has_notes_slide
                and _slide.notes_slide.notes_text_frame.text.strip()):
            continue
        _set_notes(_slide, _txt)

    out = Path(out_path) if out_path else OUT_DIR / "Module 2 - In Class Revised.pptx"
    prs.save(str(out))
    print(f"saved {out} — {len(prs.slides._sldIdLst)} slides")
    return out


if __name__ == "__main__":
    import sys as _sys
    build(_sys.argv[1] if len(_sys.argv) > 1 else None)
