"""
Add example images to the 6 revised/new slides:

  Slide 26 – AI researcher       → Anthropic logo
  Slide 31 – Rivian long-run     → Rivian R1T photo
  Slide 47 – Meta Reality Labs   → Meta Quest 3 photo  (rebuilt)
  Slide 51 – ChatGPT MC          → ChatGPT logo
  Slide 69 – AI training scale   → NVIDIA H100 photo   (rebuilt)
  Slide 72 – Alphabet scope      → Alphabet logo

For slides 47 and 69 (built fresh by the revision script) we rebuild
via clean sidecar deck with a bullets-left / picture-right layout.
For slides 26, 31, 51, 72 we add the image as an OVERLAY (no content
replacement) so the original example text stays intact.

All work via direct zip + lxml surgery on the destination deck.
Sidecars round-trip through python-pptx safely because they are
freshly built and have no NULL rels.
"""

import shutil
import sys
import zipfile
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches

HERE = Path(__file__).parent
sys.path.insert(0, str(HERE))
from _build_template_samples import (  # noqa: E402
    NAVY, GRAY, RULE, GOLD, WHITE,
    MARGIN, RULE_W, GOLD_W, SLIDE_W, SLIDE_H, FOOTER_TEXT,
    _add_text, _add_rect,
    _draw_top_bar, _draw_action_title, _draw_footer,
    _add_bulleted_list,
)
from _polish_v2 import (  # noqa: E402
    NS, Q, parse, serialize,
    REL_LAYOUT, REL_IMAGE,
    get_slide_path_at_position, next_rId,
)

TARGET = HERE / "Module 3_NEW_draft.pptx"


# ============================================================
# Sidecar build helpers (re-used from _polish_v2 conceptually)
# ============================================================

def _strip_slide(slide):
    spTree = slide.shapes._spTree
    for child in list(spTree):
        tag = child.tag.split('}', 1)[1]
        if tag in ('sp', 'grpSp', 'pic', 'graphicFrame', 'cxnSp', 'contentPart'):
            spTree.remove(child)
    slide._element.set('showMasterSp', '0')


def _new_blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _strip_slide(slide)
    return slide


# ============================================================
# Rebuild slide 47 – Meta Reality Labs with Meta Quest 3 photo
# ============================================================

def build_meta_with_image_sidecar(out_path, image_path):
    prs = Presentation()
    prs.slide_width = Emu(int(SLIDE_W))
    prs.slide_height = Emu(int(SLIDE_H))
    slide = _new_blank_slide(prs)

    _draw_top_bar(slide, "Module 3 · Part 2 · Costs")
    _draw_action_title(slide,
        "Modern sunk cost: Meta's Reality Labs has lost $50B+ since 2020")

    # Image on right (Meta Quest 3 is portrait-ish ~ 3648x5472, so render
    # vertical at about 3" wide × 4" tall)
    img_w = Inches(3.0)
    img_h = Inches(4.0)
    img_left = Inches(9.5)
    img_top = Inches(1.95)
    slide.shapes.add_picture(str(image_path), img_left, img_top,
                             width=img_w, height=img_h)
    _add_text(slide, img_left, Inches(6.0), img_w, Inches(0.35),
              "Meta Quest 3 (CC BY-SA, Wikimedia)",
              size=10, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

    # Bullets on left ~67% of slide
    bullets = [
        "Meta has poured ~$50B into Reality Labs (Metaverse, VR, AR) since 2020",
        "Wall Street keeps asking when it pays off – Zuckerberg keeps investing",
        "Past losses are sunk – ignore them when deciding what to do next",
        "Right question: does the next $10B have positive expected value?",
        "Same lesson as Waterworld – classic sunk-cost discipline, 2020s edition",
    ]
    _add_bulleted_list(
        slide,
        left=MARGIN, top=Inches(1.95),
        width=Inches(8.9), height=Inches(4.5),
        items=bullets,
        size=20, color=NAVY, bullet_color=NAVY,
        bullet_char="▪",
        bullet_size_pct=150,
        line_spacing_pts=18,
    )

    _draw_footer(slide, FOOTER_TEXT, 47)
    prs.save(out_path)


def build_ai_scale_with_image_sidecar(out_path, image_path):
    prs = Presentation()
    prs.slide_width = Emu(int(SLIDE_W))
    prs.slide_height = Emu(int(SLIDE_H))
    slide = _new_blank_slide(prs)

    _draw_top_bar(slide, "Module 3 · Part 2 · Costs")
    _draw_action_title(slide,
        "AI training: classic economies of scale at extreme cost")

    # NVIDIA H100 photo is 3840x2160 (16:9), display 5.0" wide × 2.8" tall
    img_w = Inches(5.0)
    img_h = Inches(2.8)
    img_left = Inches(7.8)
    img_top = Inches(2.3)
    slide.shapes.add_picture(str(image_path), img_left, img_top,
                             width=img_w, height=img_h)
    _add_text(slide, img_left, Inches(5.2), img_w, Inches(0.35),
              "NVIDIA H100 (CC BY-SA, Geekerwan via Wikimedia)",
              size=10, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

    bullets = [
        "Training a frontier model (GPT-5, Claude Opus 4.7) costs $500M+ in compute",
        "Once trained, marginal cost per query is fractions of a cent",
        "Spread fixed cost across 300M+ users → near-zero per-user cost",
        "NVIDIA H100/B200 GPU clusters cost $40M+ each – lumpy capital",
        "Why only a handful of labs compete in foundation models",
    ]
    _add_bulleted_list(
        slide,
        left=MARGIN, top=Inches(1.95),
        width=Inches(7.2), height=Inches(4.5),
        items=bullets,
        size=20, color=NAVY, bullet_color=NAVY,
        bullet_char="▪",
        bullet_size_pct=150,
        line_spacing_pts=18,
    )

    _draw_footer(slide, FOOTER_TEXT, 69)
    prs.save(out_path)


# ============================================================
# Replace a slide's full XML + rels via sidecar injection
# (lifted from _polish_v2 logic)
# ============================================================

def replace_slide_from_sidecar(members, sidecar_path, slide_position,
                                image_dest_filename):
    """Replace the slide at *slide_position* with the sidecar's slide 1.
    Embed the sidecar's single media file at ppt/media/<image_dest_filename>.
    """
    with zipfile.ZipFile(sidecar_path, 'r') as zf:
        slide_xml = zf.read('ppt/slides/slide1.xml')
        media_files = [n for n in zf.namelist() if n.startswith('ppt/media/')]
        assert len(media_files) == 1, f"sidecar should have 1 media, has {media_files}"
        image_bytes = zf.read(media_files[0])

    sidecar_ext = image_dest_filename.rsplit('.', 1)[1]
    # Embed image
    new_image_path = f'ppt/media/{image_dest_filename}'
    members[new_image_path] = image_bytes

    # Ensure Content_Types has Default for the extension
    ct = parse(members['[Content_Types].xml'])
    if not any(d.get('Extension') == sidecar_ext for d in ct.findall(Q('ct', 'Default'))):
        default = etree.Element(Q('ct', 'Default'))
        default.set('Extension', sidecar_ext)
        default.set('ContentType',
                    'image/jpeg' if sidecar_ext in ('jpg', 'jpeg') else f'image/{sidecar_ext}')
        ct.insert(0, default)
        members['[Content_Types].xml'] = serialize(ct)

    # Read sidecar's image rId from its rels file
    with zipfile.ZipFile(sidecar_path, 'r') as zf:
        sidecar_rels = parse(zf.read('ppt/slides/_rels/slide1.xml.rels'))
    sidecar_image_rId = None
    for rel in sidecar_rels.findall(Q('rel', 'Relationship')):
        if rel.get('Type') == REL_IMAGE:
            sidecar_image_rId = rel.get('Id')
            break

    # Replace destination slide XML + rels
    slide_path = get_slide_path_at_position(members, slide_position)
    members[slide_path] = slide_xml

    slide_dir, slide_file = slide_path.rsplit('/', 1)
    slide_rels_path = f'{slide_dir}/_rels/{slide_file}.rels'
    rels_root = etree.Element(Q('rel', 'Relationships'),
                              nsmap={None: NS['rel']})
    layout_rel = etree.SubElement(rels_root, Q('rel', 'Relationship'))
    layout_rel.set('Id', 'rId1')
    layout_rel.set('Type', REL_LAYOUT)
    layout_rel.set('Target', '../slideLayouts/slideLayout2.xml')
    image_rel = etree.SubElement(rels_root, Q('rel', 'Relationship'))
    image_rel.set('Id', sidecar_image_rId)
    image_rel.set('Type', REL_IMAGE)
    image_rel.set('Target', f'../media/{image_dest_filename}')
    members[slide_rels_path] = serialize(rels_root)


# ============================================================
# Overlay an image on an existing slide (without replacing content)
# ============================================================

def overlay_image_on_slide(members, slide_position, image_local_path,
                            image_dest_filename,
                            x_in, y_in, w_in, h_in,
                            attribution_text=None):
    """Embed an image into the package and add a <p:pic> element to the
    existing slide at *slide_position* (1-based). Position in inches.
    """
    image_bytes = image_local_path.read_bytes()
    ext = image_local_path.suffix.lstrip('.').lower()
    new_image_path = f'ppt/media/{image_dest_filename}'
    members[new_image_path] = image_bytes

    # Content_Types Default
    ct = parse(members['[Content_Types].xml'])
    if not any(d.get('Extension') == ext for d in ct.findall(Q('ct', 'Default'))):
        default = etree.Element(Q('ct', 'Default'))
        default.set('Extension', ext)
        default.set('ContentType',
                    'image/jpeg' if ext in ('jpg', 'jpeg') else f'image/{ext}')
        ct.insert(0, default)
        members['[Content_Types].xml'] = serialize(ct)

    # Slide path + rels
    slide_path = get_slide_path_at_position(members, slide_position)
    slide_dir, slide_file = slide_path.rsplit('/', 1)
    slide_rels_path = f'{slide_dir}/_rels/{slide_file}.rels'

    # Add image rel
    slide_rels = parse(members[slide_rels_path])
    img_rId = next_rId(slide_rels)
    img_rel = etree.SubElement(slide_rels, Q('rel', 'Relationship'))
    img_rel.set('Id', img_rId)
    img_rel.set('Type', REL_IMAGE)
    img_rel.set('Target', f'../media/{image_dest_filename}')
    members[slide_rels_path] = serialize(slide_rels)

    # Build the <p:pic> element XML and append to spTree
    EMU = 914400
    x = int(x_in * EMU)
    y = int(y_in * EMU)
    w = int(w_in * EMU)
    h = int(h_in * EMU)
    pic_xml = (
        f'<p:pic xmlns:p="{NS["p"]}" xmlns:a="{NS["a"]}" xmlns:r="{NS["r"]}">'
        f'<p:nvPicPr>'
          f'<p:cNvPr id="9999" name="ExampleImage"/>'
          f'<p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>'
          f'<p:nvPr/>'
        f'</p:nvPicPr>'
        f'<p:blipFill>'
          f'<a:blip r:embed="{img_rId}"/>'
          f'<a:stretch><a:fillRect/></a:stretch>'
        f'</p:blipFill>'
        f'<p:spPr>'
          f'<a:xfrm><a:off x="{x}" y="{y}"/><a:ext cx="{w}" cy="{h}"/></a:xfrm>'
          f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
        f'</p:spPr>'
        f'</p:pic>'
    )
    pic_el = etree.fromstring(pic_xml)

    slide = parse(members[slide_path])
    spTree = slide.find(Q('p', 'cSld') + '/' + Q('p', 'spTree'))
    spTree.append(pic_el)

    # Optional attribution caption
    if attribution_text:
        cap_xml = (
            f'<p:sp xmlns:p="{NS["p"]}" xmlns:a="{NS["a"]}">'
            f'<p:nvSpPr>'
              f'<p:cNvPr id="9998" name="AttrText"/>'
              f'<p:cNvSpPr txBox="1"/><p:nvPr/>'
            f'</p:nvSpPr>'
            f'<p:spPr>'
              f'<a:xfrm><a:off x="{x}" y="{y + h + int(0.05 * EMU)}"/>'
              f'<a:ext cx="{w}" cy="{int(0.3 * EMU)}"/></a:xfrm>'
              f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
              f'<a:noFill/>'
            f'</p:spPr>'
            f'<p:txBody>'
              f'<a:bodyPr wrap="square" lIns="0" rIns="0" tIns="0" bIns="0"/>'
              f'<a:lstStyle/>'
              f'<a:p><a:pPr algn="ctr"/>'
              f'<a:r><a:rPr sz="900" i="1">'
              f'<a:solidFill><a:srgbClr val="555B66"/></a:solidFill>'
              f'<a:latin typeface="Calibri"/></a:rPr>'
              f'<a:t>{attribution_text}</a:t></a:r></a:p>'
            f'</p:txBody>'
            f'</p:sp>'
        )
        spTree.append(etree.fromstring(cap_xml))

    members[slide_path] = serialize(slide)


# ============================================================
# Main
# ============================================================

SIDECAR_META = HERE / "_sidecar_meta_v2.pptx"
SIDECAR_AI = HERE / "_sidecar_ai_v2.pptx"


def main():
    print(f"Reading {TARGET.name}")
    with zipfile.ZipFile(TARGET, 'r') as zf:
        members = {n: zf.read(n) for n in zf.namelist()}

    # 1) Rebuild slide 47 (Meta Reality Labs) with Quest 3 image
    print("[1] Rebuild slide 47 (Meta Reality Labs) with Quest 3 image")
    build_meta_with_image_sidecar(SIDECAR_META, HERE / "_meta_quest.jpg")
    replace_slide_from_sidecar(members, SIDECAR_META, 47, "_meta_quest.jpg")

    # 2) Rebuild slide 69 (AI training) with NVIDIA H100 image
    print("[2] Rebuild slide 69 (AI training scale) with NVIDIA H100 image")
    build_ai_scale_with_image_sidecar(SIDECAR_AI, HERE / "_nvidia_h100.png")
    replace_slide_from_sidecar(members, SIDECAR_AI, 69, "_nvidia_h100.png")

    # 3) Overlay Anthropic logo on slide 26
    print("[3] Overlay Anthropic logo on slide 26")
    overlay_image_on_slide(
        members, 26, HERE / "_anthropic_logo.png", "_anthropic_logo.png",
        x_in=9.6, y_in=5.4, w_in=3.2, h_in=0.36,
        attribution_text=None,
    )

    # 4) Overlay Rivian R1T photo on slide 31
    print("[4] Overlay Rivian R1T photo on slide 31")
    overlay_image_on_slide(
        members, 31, HERE / "_rivian.jpg", "_rivian.jpg",
        x_in=8.0, y_in=4.5, w_in=4.8, h_in=2.4,
        attribution_text="Rivian R1T (CC BY-SA, Wikimedia)",
    )

    # 5) Overlay ChatGPT logo on slide 51
    print("[5] Overlay ChatGPT logo on slide 51")
    overlay_image_on_slide(
        members, 51, HERE / "_chatgpt_logo.png", "_chatgpt_logo.png",
        x_in=10.3, y_in=4.5, w_in=2.5, h_in=2.5,
        attribution_text=None,
    )

    # 6) Overlay Alphabet logo on slide 72
    print("[6] Overlay Alphabet logo on slide 72")
    overlay_image_on_slide(
        members, 72, HERE / "_alphabet_logo.png", "_alphabet_logo.png",
        x_in=9.5, y_in=5.7, w_in=3.3, h_in=0.55,
        attribution_text=None,
    )

    # Write atomically
    tmp = TARGET.with_suffix('.pptx.tmp')
    with zipfile.ZipFile(tmp, 'w', zipfile.ZIP_DEFLATED) as zf:
        for n, d in members.items():
            zf.writestr(n, d)
    tmp.replace(TARGET)
    print(f"\nWrote {TARGET.name}")


if __name__ == "__main__":
    main()
