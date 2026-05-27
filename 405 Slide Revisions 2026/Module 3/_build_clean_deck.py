"""
Build a clean Module 3 deck from scratch, using ONLY the six template layouts
defined in `_build_template_samples.py`.

Goal: every slide in this deck uses one of six layout types (title, section
header, content bulleted, content two-column, poll, closing synthesis), all
on the Blank layout, so PowerPoint's Layout dropdown stays clean.

Build is by batches – front matter (1-6), then §1.1 Short Run (7-22), etc.

Output: `Module 3_clean.pptx`
"""

import copy
import math
import re
import shutil
import zipfile
from pathlib import Path

from lxml import etree as ET
from pptx import Presentation
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Cm, Inches, Pt

# Reuse all primitives from the template script (single source of truth).
from _build_template_samples import (
    FADED,
    FOOTER_TEXT,
    GOLD,
    GOLD_W,
    GRAY,
    MARGIN,
    MODULE_AGENDA,
    NAVY,
    RULE,
    RULE_W,
    SLIDE_H,
    SLIDE_W,
    WHITE,
    _add_bulleted_list,
    _add_rect,
    _add_text,
    _blank_slide,
    _draw_action_title,
    _set_bullet_char,
)

OUT_DIR = Path(__file__).parent


# --------------------------------------------------------------------------
# Title-case top bar – replaces the all-caps default from the template script.
# The user prefers academic-paper title case (each major word capitalized,
# small connectors like "and", "of", "for", "the" stay lowercase).
# --------------------------------------------------------------------------

def _draw_top_bar_tc(slide, section_tag):
    """Navy top bar with title-cased section tag (no uppercase forcing)."""
    bar_h = Inches(0.42)
    _add_rect(slide, 0, 0, SLIDE_W, bar_h, NAVY)
    _add_text(slide, MARGIN, 0, Inches(12), bar_h,
              section_tag, size=16, bold=True,
              color=WHITE, font="Calibri",
              anchor=MSO_ANCHOR.MIDDLE)


def _draw_footer(slide, footer_text, page_num):
    """Footer rule + gold accent + footer text/page number, with larger type
    than the template default so handout printouts remain legible."""
    _add_rect(slide, 0, Inches(7.15), SLIDE_W, Inches(0.02), RULE)
    _add_rect(slide, MARGIN, Inches(7.135), GOLD_W, Inches(0.05), GOLD)
    _add_text(slide, MARGIN, Inches(7.20), Inches(11), Inches(0.32),
              footer_text, size=12, color=GRAY)
    _add_text(slide, Inches(12.5), Inches(7.20), Inches(0.6), Inches(0.32),
              str(page_num), size=12, color=GRAY, align=PP_ALIGN.RIGHT)


# --------------------------------------------------------------------------
# Speaker-notes helper
# --------------------------------------------------------------------------

def _set_notes(slide, text):
    """Replace the slide's speaker notes with *text*."""
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.clear()
    notes_tf.text = text


# --------------------------------------------------------------------------
# Reusable shape primitives for diagrams
# --------------------------------------------------------------------------

def _add_filled_box(slide, left, top, width, height, label, *,
                    fill=NAVY, text_color=WHITE, line=None,
                    size=18, bold=True, font="Calibri"):
    """Filled rectangle with centered text."""
    left, top, width, height = int(left), int(top), int(width), int(height)
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height,
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = text_color
    return shp


def _add_outlined_box(slide, left, top, width, height, label, *,
                      line=NAVY, text_color=NAVY, fill=WHITE,
                      size=18, bold=True, line_w=1.25, font="Calibri"):
    """Outlined rectangle (white fill) with centered text."""
    left, top, width, height = int(left), int(top), int(width), int(height)
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height,
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = text_color
    return shp


def _add_convention_box(slide, left, top, width, height, *,
                          prefix=None, body=None, runs=None,
                          fill_rgb=None, border=None, line_w=1.0,
                          corner_pct=0.12, size=15, align=PP_ALIGN.LEFT,
                          font="Calibri", pad_h=None, pad_v=None,
                          line_spacing_pct=None):
    """Cream-fill / navy-border rounded-rect explanation callout.

    The "Convention" textbox pattern from slide 14 generalised — use it
    anywhere a slide needs a compact, visually-distinct box for a
    short conceptual explanation or notational convention.  Sits well
    below a table, beside a hero formula, or as a slide-wide footer.

    Two ways to populate the text:
      • ``prefix`` (bold) + ``body`` (regular) — simplest path; matches
        slide 14's "Convention:  <text>" pattern.
      • ``runs`` — a list of ``(text, {"bold": .., "italic": .., ...})``
        tuples for finer-grained styling (multi-line, mixed formatting).

    Style defaults follow the course-layer CLAUDE.md "Convention callout
    box" spec — cream fill, thin primary-color border, slight rounding,
    primary-color text.  Override ``fill_rgb`` / ``border`` only when you
    need a different accent.
    """
    fill = fill_rgb if fill_rgb is not None else RGBColor(0xFD, 0xF6, 0xE6)
    border = border if border is not None else NAVY

    left, top, width, height = int(left), int(top), int(width), int(height)
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border
    shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    try:
        shp.adjustments[0] = corner_pct
    except Exception:
        pass

    # Inset text box so the rounded corners breathe — matches slide 14.
    pad_h = Inches(0.20) if pad_h is None else pad_h
    pad_v = Inches(0.12) if pad_v is None else pad_v
    tb = slide.shapes.add_textbox(
        left + int(pad_h), top + int(pad_v),
        width - 2 * int(pad_h), height - 2 * int(pad_v),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0); tf.margin_bottom = Inches(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    def _style_run(r, opts):
        r.font.name = opts.get('font', font)
        r.font.size = Pt(opts.get('size', size))
        r.font.bold = opts.get('bold', False)
        r.font.italic = opts.get('italic', False)
        r.font.color.rgb = opts.get('color', NAVY)

    if runs is not None:
        first = True
        for entry in runs:
            text, opts = entry if isinstance(entry, tuple) else (entry, {})
            if opts.get('newline') and not first:
                p = tf.add_paragraph()
            elif first:
                p = tf.paragraphs[0]
            else:
                # Same paragraph — append run to the most-recent paragraph.
                p = tf.paragraphs[-1]
            p.alignment = align
            r = p.add_run()
            r.text = text
            _style_run(r, opts)
            first = False
    else:
        p = tf.paragraphs[0]
        p.alignment = align
        if prefix:
            r1 = p.add_run(); r1.text = prefix
            _style_run(r1, {'bold': True, 'color': NAVY, 'size': size})
        if body:
            r2 = p.add_run(); r2.text = body
            _style_run(r2, {'color': NAVY, 'size': size})

    if line_spacing_pct is not None:
        for p_obj in tf.paragraphs:
            pPr = p_obj._p.get_or_add_pPr()
            for old in pPr.findall(qn('a:lnSpc')):
                pPr.remove(old)
            lnSpc = ET.Element(qn('a:lnSpc'))
            spcPct = ET.SubElement(lnSpc, qn('a:spcPct'))
            spcPct.set('val', str(int(line_spacing_pct * 1000)))
            pPr.insert(0, lnSpc)
    return shp


def _add_rounded_filled_box(slide, left, top, width, height, label, *,
                             fill=NAVY, text_color=WHITE, line=None,
                             size=18, bold=True, font="Calibri",
                             corner_pct=0.06, shadow=True):
    """Rounded-corner filled rectangle with centered text and soft drop shadow.

    Mirrors :func:`_add_filled_box` but renders ``MSO_SHAPE.ROUNDED_RECTANGLE``
    with the corner-adjust set to ``corner_pct`` (6 % per course CLAUDE.md
    "slight rounding") and a soft drop shadow via :func:`_add_drop_shadow`.
    """
    left, top, width, height = int(left), int(top), int(width), int(height)
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
    )
    try:
        shp.adjustments[0] = corner_pct
    except Exception:
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    if shadow:
        _add_drop_shadow(shp)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = text_color
    return shp


def _add_arrow(slide, start_xy, end_xy, *, color=NAVY, weight_pt=1.5,
               head=True, dash=None, head_size='med'):
    """Draw a line/arrow from start to end (in EMU/Inches values).

    EMU coordinates MUST be integers — PowerPoint rejects decimal values
    in <a:off>/<a:ext> and refuses to open the file. Cast to int defensively.

    ``dash`` accepts any OOXML preset-dash name (e.g., ``"dash"``,
    ``"dashDot"``, ``"sysDash"``).  Default ``None`` = solid line.

    ``head_size`` is the OOXML preset arrowhead size — one of ``'sm'``,
    ``'med'`` (default), or ``'lg'``.  Width and height are set
    together, so passing ``'lg'`` gives a noticeably larger tip while
    leaving the line weight unchanged.
    """
    sx, sy = int(start_xy[0]), int(start_xy[1])
    ex, ey = int(end_xy[0]), int(end_xy[1])
    line = slide.shapes.add_connector(1, sx, sy, ex, ey)  # 1 = STRAIGHT
    line.line.color.rgb = color
    line.line.width = Pt(weight_pt)
    ln = line.line._get_or_add_ln()
    if dash is not None:
        for old in ln.findall(qn('a:prstDash')):
            ln.remove(old)
        prst = ET.SubElement(ln, qn('a:prstDash'))
        prst.set('val', dash)
    if head:
        tailEnd = ET.SubElement(ln, qn('a:tailEnd'))
        tailEnd.set('type', 'triangle')
        tailEnd.set('w', head_size)
        tailEnd.set('h', head_size)
    return line


def _add_wavy_line(slide, x_start, x_end, y_center, *,
                    amplitude=None, cycles=1.75, segments=36,
                    color=NAVY, weight_pt=1.5):
    """Horizontal sinusoidal line from x_start to x_end at y_center.

    Renders a polyline approximation of ``sin(2π · t · cycles)`` inside a
    custGeom shape so the line reads as a gentle wave rather than a
    straight connector.  ``amplitude`` is the peak-to-baseline height in
    EMU; defaults to ~0.04".  ``segments`` controls how finely the wave
    is discretised — 30+ is smooth enough that the polyline reads as a
    curve.  No arrowhead.
    """
    if amplitude is None:
        amplitude = Inches(0.04)
    L = int(x_end - x_start)
    A = int(amplitude)
    if L == 0:
        return None
    flip = "1" if L < 0 else "0"
    bbox_left = int(min(x_start, x_end))
    bbox_top = int(y_center - A)
    bbox_w = abs(L)
    bbox_h = 2 * A

    pts = []
    for i in range(segments + 1):
        t = i / segments
        lx = int(round(t * 100000))
        sin_val = math.sin(2 * math.pi * t * cycles)
        # Path coords: y=0 is top.  Centre at 50000; +1·amplitude → top
        # (0), −1·amplitude → bottom (100000).
        ly = int(round(50000 - sin_val * 50000))
        pts.append((lx, ly))

    path_segs = [f'<a:moveTo><a:pt x="{pts[0][0]}" y="{pts[0][1]}"/></a:moveTo>']
    for lx, ly in pts[1:]:
        path_segs.append(f'<a:lnTo><a:pt x="{lx}" y="{ly}"/></a:lnTo>')
    path_inner = ''.join(path_segs)

    color_hex = f'{color[0]:02X}{color[1]:02X}{color[2]:02X}'
    width_emu = int(weight_pt * 12700)

    P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    A_NS_LOCAL = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    sp_xml = (
        f'<p:sp xmlns:p="{P_NS}" xmlns:a="{A_NS_LOCAL}">'
        f'<p:nvSpPr><p:cNvPr id="0" name="WavyLine"/>'
        f'<p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr>'
        f'<a:xfrm flipH="{flip}">'
        f'<a:off x="{bbox_left}" y="{bbox_top}"/>'
        f'<a:ext cx="{bbox_w}" cy="{bbox_h}"/>'
        f'</a:xfrm>'
        f'<a:custGeom>'
        f'<a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/>'
        f'<a:rect l="0" t="0" r="0" b="0"/>'
        f'<a:pathLst>'
        f'<a:path w="100000" h="100000" fill="none">'
        f'{path_inner}'
        f'</a:path>'
        f'</a:pathLst>'
        f'</a:custGeom>'
        f'<a:noFill/>'
        f'<a:ln w="{width_emu}" cap="rnd">'
        f'<a:solidFill><a:srgbClr val="{color_hex}"/></a:solidFill>'
        f'</a:ln>'
        f'</p:spPr>'
        f'</p:sp>'
    )
    elem = ET.fromstring(sp_xml)
    slide.shapes._spTree.append(elem)
    return elem


def _add_arrow_shape(slide, left, top, width, height, *,
                     direction="right", fill=GOLD, line=None):
    """Block arrow shape (the 'we-are-here' indicator).

    direction: "right" (default), "left", "up", or "down".
    """
    left, top, width, height = int(left), int(top), int(width), int(height)
    geom_map = {
        "left": MSO_SHAPE.LEFT_ARROW,
        "right": MSO_SHAPE.RIGHT_ARROW,
        "up": MSO_SHAPE.UP_ARROW,
        "down": MSO_SHAPE.DOWN_ARROW,
    }
    geom = geom_map.get(direction, MSO_SHAPE.RIGHT_ARROW)
    shp = slide.shapes.add_shape(geom, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


# --------------------------------------------------------------------------
# Layout 1 — Title slide
# --------------------------------------------------------------------------

def make_title_slide(prs):
    slide = _blank_slide(prs)
    _add_text(slide, MARGIN, Inches(2.0), RULE_W, Inches(1.3),
              "Production and Costs",
              size=60, bold=True, color=NAVY, font="Calibri",
              align=PP_ALIGN.CENTER)
    _add_text(slide, MARGIN, Inches(3.35), RULE_W, Inches(0.75),
              "Module 3",
              size=40, bold=True, color=GOLD, font="Calibri",
              align=PP_ALIGN.CENTER)
    accent_w = Inches(4.0)
    accent_x = (SLIDE_W - accent_w) // 2
    _add_rect(slide, accent_x, Inches(4.4), accent_w, Inches(0.06), GOLD)
    _add_text(slide, MARGIN, Inches(4.8), RULE_W, Inches(0.55),
              "Management 405  ·  EMBA",
              size=26, bold=True, color=GRAY, font="Calibri",
              align=PP_ALIGN.CENTER)
    _add_text(slide, MARGIN, Inches(5.5), RULE_W, Inches(0.5),
              "Prof. Nico Voigtlaender  ·  UCLA Anderson",
              size=22, color=GRAY, font="Calibri",
              align=PP_ALIGN.CENTER)
    _add_rect(slide, 0, Inches(7.15), SLIDE_W, Inches(0.02), RULE)
    _add_rect(slide, MARGIN, Inches(7.135), GOLD_W, Inches(0.05), GOLD)
    return slide


# --------------------------------------------------------------------------
# Layout 2 — Section Header / Agenda (parameterized)
# --------------------------------------------------------------------------

def make_section_agenda(prs, page_num, *, current_part_idx=None,
                        current_sub_idx=None,
                        section_tag="Module 3 · Agenda",
                        title="Agenda"):
    """Render an agenda / section-divider slide.

    ``current_part_idx`` makes that Part navy and fades the others.
    ``current_sub_idx`` (optional, only meaningful with current_part_idx)
    further restricts the navy highlight inside the current Part to a
    single sub-bullet — the other subs render in faded gray, alongside
    Part 2.  Used by intra-Part dividers (e.g., the new Short Run agenda
    on page 13, or the Long Run agenda on page 31).
    """
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, section_tag)
    _draw_action_title(slide, title)

    y = Inches(1.85)
    part_title_h = Inches(0.6)
    sub_list_h = Inches(1.35)
    block_gap = Inches(0.15)

    for idx, part in enumerate(MODULE_AGENDA):
        if current_part_idx is None:
            color = NAVY  # preview: highlight all Parts
        else:
            color = NAVY if idx == current_part_idx else FADED

        # 2026-05-19: per-sub fade.  When this is the current Part AND a
        # specific sub-index is highlighted, give every other sub the
        # FADED color so the deck shows "you are HERE" within the Part.
        if (current_part_idx is not None and idx == current_part_idx
                and current_sub_idx is not None):
            sub_colors = [
                NAVY if i == current_sub_idx else FADED
                for i in range(len(part["subs"]))
            ]
        else:
            sub_colors = None

        _add_text(slide, MARGIN, y, RULE_W, part_title_h,
                  part["title"], size=30, bold=True, color=color,
                  font="Calibri")
        y += part_title_h

        _add_bulleted_list(
            slide,
            left=MARGIN + Inches(0.4),
            top=y,
            width=RULE_W - Inches(0.4),
            height=sub_list_h,
            items=part["subs"],
            size=24, color=color, bullet_color=color,
            # 2026-05-19: sub-bullets tightened from 10 → 3 pt to match
            # the deck-wide rhythm (sub-bullets cluster under their
            # parent rather than breathing out).  Applied to both
            # Parts on every agenda divider in the deck.
            line_spacing_pts=3,
            autonum_scheme='alphaLcPeriod',
            colors=sub_colors,
        )
        y += sub_list_h + block_gap

    _draw_footer(slide, FOOTER_TEXT, page_num)
    return slide


# --------------------------------------------------------------------------
# Layout 3 — Content bulleted
# --------------------------------------------------------------------------

def make_content_bulleted(prs, page_num, section_tag, title, bullets, *,
                          size=24, sub_size=None, line_spacing_pts=18,
                          sub_line_spacing_pts=None,
                          extras=None, bullets_top=None):
    """bullets: list of (text, level) tuples OR plain strings (level=0).

    ``bullets_top`` overrides the default body-region start (Inches(1.85))
    — useful when the slide also hosts large diagrams below the bullets
    and the bullets need to lift up to avoid overlap.

    ``sub_line_spacing_pts`` overrides the legacy
    ``max(6, line_spacing_pts - 8)`` formula for sub-bullet space-before.
    """
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, section_tag)
    _draw_action_title(slide, title)

    normalized = [(b, 0) if isinstance(b, str) else b for b in bullets]

    if bullets_top is None:
        bullets_top = Inches(1.85)

    _add_hierarchical_bullets(
        slide,
        left=MARGIN,
        top=bullets_top,
        width=RULE_W,
        height=Inches(5.0),
        items=normalized,
        size=size,
        sub_size=sub_size,
        line_spacing_pts=line_spacing_pts,
        sub_line_spacing_pts=sub_line_spacing_pts,
    )

    if extras is not None:
        extras(slide)

    _draw_footer(slide, FOOTER_TEXT, page_num)
    return slide


def _inject_bullet_lst_style(tf, *, size, sub_size,
                              main_color=NAVY, sub_color=GRAY,
                              main_char='▪', sub_char='–',
                              main_space_pts=None, sub_space_pts=None):
    """Inject <a:lstStyle> into a text frame defining per-level defaults.

    Defines lvl1pPr (main bullets, lvl="0") and lvl2pPr (sub bullets,
    lvl="1") with bullet character, indent, color, font, and space-
    before defaults.  Enables PowerPoint's native Tab / Shift+Tab to
    auto-reformat a bullet when its level changes: runs / paragraphs
    that do NOT explicitly override these properties inherit them.

    main_space_pts / sub_space_pts: if set, the <a:spcBef> for each
    level.  Sub-bullets default to a tighter spacing than main bullets
    (3 pt vs ~10 pt) per the deck's pedagogical preference: sub-bullets
    visually cluster under their parent, main bullets give the eye a
    larger break.
    """
    txBody = tf._txBody
    # Remove any prior lstStyle
    for old in txBody.findall(qn('a:lstStyle')):
        txBody.remove(old)
    lstStyle = ET.Element(qn('a:lstStyle'))
    levels = [
        ('a:lvl1pPr', 342900, -342900, main_color, main_char, size, main_space_pts),
        ('a:lvl2pPr', 685800, -228600, sub_color, sub_char, sub_size, sub_space_pts),
    ]
    for tag, mar_l, indent, color, char, sz, space_pts in levels:
        lvl = ET.SubElement(lstStyle, qn(tag))
        lvl.set('marL', str(mar_l))
        lvl.set('indent', str(indent))
        # spcBef must precede bullet attributes (OOXML schema order).
        if space_pts is not None:
            spcBef = ET.SubElement(lvl, qn('a:spcBef'))
            spcPts = ET.SubElement(spcBef, qn('a:spcPts'))
            spcPts.set('val', str(int(space_pts * 100)))
        bc = ET.SubElement(lvl, qn('a:buClr'))
        sc = ET.SubElement(bc, qn('a:srgbClr'))
        sc.set('val', '{:02X}{:02X}{:02X}'.format(color[0], color[1], color[2]))
        bf = ET.SubElement(lvl, qn('a:buFont'))
        bf.set('typeface', 'Calibri')
        bch = ET.SubElement(lvl, qn('a:buChar'))
        bch.set('char', char)
        drp = ET.SubElement(lvl, qn('a:defRPr'))
        drp.set('sz', str(int(sz * 100)))
        drp.set('b', '0')
        sf = ET.SubElement(drp, qn('a:solidFill'))
        sfc = ET.SubElement(sf, qn('a:srgbClr'))
        sfc.set('val', '{:02X}{:02X}{:02X}'.format(color[0], color[1], color[2]))
        lt = ET.SubElement(drp, qn('a:latin'))
        lt.set('typeface', 'Calibri')
    # Insert lstStyle after bodyPr (proper element order in <a:txBody>)
    bodyPr = txBody.find(qn('a:bodyPr'))
    if bodyPr is not None:
        bodyPr.addnext(lstStyle)
    else:
        txBody.insert(0, lstStyle)


def _add_hierarchical_bullets(slide, left, top, width, height, items,
                              *, size=24, sub_size=None, line_spacing_pts=18,
                              sub_line_spacing_pts=None):
    """Render bullets with indent levels.

    Bullet item forms:
        (text, level)                       — simple, defaults from level
        (text, level, opts)                 — opts dict overrides
        (runs_list, level, opts)            — multi-run paragraph

    text:
        - str  → a single run with text
        - ''   → empty paragraph (visual spacer; no run)
        - list → multi-run: list of ``(run_text, run_opts)`` tuples

    Paragraph-level opts (all optional):
        bullet_style: 'main' (▪ NAVY) | 'sub' (– GRAY) | 'arrow' (no bullet,
            plain left-indent; the run text supplies the leader char like
            "→" or Wingdings) | 'none'
        mar_l, indent: bullet positioning (EMU)
        space_before_pts: spcBef in pts (overrides legacy formula)
        size, color, bold, italic: defaults applied to every run unless
            the run_opts override them.

    Run-level opts (run_opts in a runs_list tuple):
        font_name (default 'Calibri'), size, color, bold, italic,
        underline (bool), wingdings (bool — emits <a:sym typeface="Wingdings"/>
        so a private-use-area character renders as its Wingdings glyph).

    Each paragraph also receives a ``lvl="N"`` attribute when level > 0
    so PowerPoint's Tab / Shift-Tab outline navigation can find an
    explicit outline level.
    """
    if sub_size is None:
        sub_size = size - 4

    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    # 2026-05-18: inject <a:lstStyle> with per-level defaults so that
    # PowerPoint's Tab / Shift+Tab keyboard shortcuts auto-reformat the
    # bullet when its outline level changes.  Default 'main'/'sub'
    # bullets (no custom indent, level < 2) skip explicit pPr/run
    # styling — they inherit from lstStyle.  Bullets with custom indent
    # OR level >= 2 OR explicit size/color overrides still apply
    # explicit styling.
    # 2026-05-18 (later): also moved space-before into lstStyle so a
    # paragraph demoted via Tab picks up the sub-bullet's tighter spcBef
    # automatically.  Sub-bullet default: 3 pt (was 6 pt) — per user
    # preference, sub-bullets cluster tightly under their parent.
    sub_default_space = (sub_line_spacing_pts if sub_line_spacing_pts is not None
                          else 3)
    _inject_bullet_lst_style(tf, size=size, sub_size=sub_size,
                              main_space_pts=line_spacing_pts,
                              sub_space_pts=sub_default_space)

    for i, item in enumerate(items):
        if len(item) == 2:
            text, level = item
            opts = {}
        else:
            text, level, opts = item
            opts = opts or {}

        # Determine inheritance up-front so spcBef logic below can use it.
        style = opts.get('bullet_style', 'main' if level == 0 else 'sub')
        has_custom_indent = 'mar_l' in opts or 'indent' in opts
        inherits_from_lst = (style in ('main', 'sub')
                              and not has_custom_indent
                              and level < 2)

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        pPr = p._p.get_or_add_pPr()

        # Space-before.  When the paragraph inherits from lstStyle and the
        # caller has not overridden, skip the per-pPr spcBef so that Tab-
        # induced level changes pick up the new level's spcBef from
        # lstStyle.  Otherwise set explicitly.
        if i > 0:
            sp_override = opts.get('space_before_pts')
            if sp_override is not None:
                spcBef = ET.SubElement(pPr, qn('a:spcBef'))
                pts = ET.SubElement(spcBef, qn('a:spcPts'))
                pts.set('val', str(sp_override * 100))
            elif not inherits_from_lst:
                if level == 0:
                    sp = line_spacing_pts
                else:
                    sp = sub_default_space
                spcBef = ET.SubElement(pPr, qn('a:spcBef'))
                pts = ET.SubElement(spcBef, qn('a:spcPts'))
                pts.set('val', str(sp * 100))
            # else: inherit from lstStyle.

        # Outline-level attribute (enables PowerPoint Tab/Shift-Tab)
        if level > 0:
            pPr.set('lvl', str(level))

        # Bullet styling — default 'main' (lvl 0) and 'sub' (lvl 1)
        # inherit from lstStyle when no custom indent is requested.
        # 'arrow' and 'none' suppress the lstStyle bullet via <a:buNone/>.
        # Level >= 2 still uses explicit _set_bullet_char (lstStyle here
        # only defines lvl1pPr and lvl2pPr).
        if style == 'main':
            if not inherits_from_lst:
                _set_bullet_char(p, char='▪', color=NAVY,
                                  mar_l=opts.get('mar_l', 342900),
                                  indent=opts.get('indent', -342900),
                                  size_pct=100)
        elif style == 'sub':
            if not inherits_from_lst:
                default_mar = 342900 + level * 342900
                _set_bullet_char(p, char='–', color=GRAY,
                                  mar_l=opts.get('mar_l', default_mar),
                                  indent=opts.get('indent', -228600),
                                  size_pct=100)
        elif style == 'arrow':
            # Plain left-indent, no bullet glyph; user text supplies leader.
            pPr.set('marL', str(opts.get('mar_l', 457200)))
            if 'indent' in opts:
                pPr.set('indent', str(opts['indent']))
            # Explicit <a:buNone/> so the lstStyle bullet doesn't bleed
            # through.
            ET.SubElement(pPr, qn('a:buNone'))
        elif style == 'none':
            ET.SubElement(pPr, qn('a:buNone'))

        # Empty paragraph (spacer) — no run; suppress the bullet so the
        # empty line doesn't render an orphan glyph.
        if text == '':
            if inherits_from_lst:
                ET.SubElement(pPr, qn('a:buNone'))
            continue

        # Normalize text to a runs list
        if isinstance(text, str):
            runs = [(text, {})]
        else:
            runs = text  # already a list of (run_text, run_opts) tuples

        # Paragraph-level explicit overrides (None = inherit from lstStyle).
        para_size_override = opts.get('size')
        para_color_override = opts.get('color')
        para_bold_override = opts.get('bold')
        para_italic_override = opts.get('italic')

        for run_text, run_opts in runs:
            run_opts = run_opts or {}

            # Inline OMML math zone — append <a14:m><m:oMath>...</m:oMath></a14:m>
            # as a sibling of the regular <a:r> runs so the formula renders
            # in-line with surrounding prose.  ``run_text`` is interpreted
            # as raw OMML content (e.g., from `_formula_mp_ratio` or the
            # `_omml_*` builders).
            if run_opts.get('omml'):
                om_sz = run_opts.get('size')
                if om_sz is None:
                    om_sz = para_size_override
                if om_sz is None:
                    om_sz = size if level == 0 else sub_size
                om_clr = run_opts.get('color') or para_color_override
                if om_clr is None:
                    om_clr = NAVY if level == 0 else GRAY
                sz_centi = int(om_sz * 100)
                clr_hex = '{:02X}{:02X}{:02X}'.format(
                    om_clr[0], om_clr[1], om_clr[2])
                a14_xml = (
                    f'<a14:m xmlns:a14="{A14_NS}" '
                    f'xmlns:m="{M_NS}" xmlns:a="{A_NS}">'
                    f'<m:oMath>{run_text}</m:oMath>'
                    f'</a14:m>'
                )
                a14_elem = ET.fromstring(a14_xml)
                for r_elem in a14_elem.iter(qn('m:r')):
                    arPr = r_elem.find(qn('a:rPr'))
                    if arPr is None:
                        arPr = ET.Element(qn('a:rPr'))
                        r_elem.insert(0, arPr)
                    arPr.set('sz', str(sz_centi))
                    if arPr.get('lang') is None:
                        arPr.set('lang', 'en-US')
                    for sf in arPr.findall(qn('a:solidFill')):
                        arPr.remove(sf)
                    sf = ET.SubElement(arPr, qn('a:solidFill'))
                    srgb = ET.SubElement(sf, qn('a:srgbClr'))
                    srgb.set('val', clr_hex)
                p._p.append(a14_elem)
                continue

            run = p.add_run()
            run.text = run_text
            run.font.name = run_opts.get('font_name', 'Calibri')

            # Size: run-opt > para-opt > inherit (for inherits_from_lst
            # cases) or fall back to level default (for explicit-styled
            # paragraphs).
            run_sz = run_opts.get('size')
            sz = run_sz if run_sz is not None else para_size_override
            if sz is None and not inherits_from_lst:
                sz = size if level == 0 else sub_size
            if sz is not None:
                run.font.size = Pt(sz)

            run_clr = run_opts.get('color')
            clr = run_clr if run_clr is not None else para_color_override
            if clr is None and not inherits_from_lst:
                clr = NAVY if level == 0 else GRAY
            if clr is not None:
                run.font.color.rgb = clr

            # Bold: only set if explicitly requested.  Default = inherit
            # (lstStyle defRPr has b="0").
            run_b = run_opts.get('bold')
            b = run_b if run_b is not None else para_bold_override
            if b is not None:
                run.font.bold = b

            it = run_opts.get('italic', para_italic_override)
            if it is not None:
                run.font.italic = it
            if run_opts.get('underline'):
                run.font.underline = True
            if run_opts.get('wingdings'):
                # Add <a:sym typeface="Wingdings"/> so private-use-area
                # characters render as their Wingdings glyphs.
                rPr = run._r.find(qn('a:rPr'))
                if rPr is None:
                    rPr = run._r.makeelement(qn('a:rPr'), {})
                    run._r.insert(0, rPr)
                sym = ET.SubElement(rPr, qn('a:sym'))
                sym.set('typeface', 'Wingdings')

    return box


# --------------------------------------------------------------------------
# Layout-3 variant — diagram canvas (action title + free-form shapes below).
# Used for the agenda flowchart and the Big Picture diagram so they live
# inside the same visual chrome as content slides but render a diagram
# instead of bullets.
# --------------------------------------------------------------------------

def make_diagram_slide(prs, page_num, section_tag, title, draw_diagram):
    """Action-title slide with free-form diagram region below.

    draw_diagram: callable(slide) that renders the diagram in the body
    region (approximately y = 1.85 to 6.95).
    """
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, section_tag)
    _draw_action_title(slide, title)
    draw_diagram(slide)
    _draw_footer(slide, FOOTER_TEXT, page_num)
    return slide


# --------------------------------------------------------------------------
# Layout 5 — Poll slide (A./B./C./D. options + POLL pill, no QR box)
# --------------------------------------------------------------------------

def _draw_poll_pill(slide, *, position='top-right',
                     fill=NAVY, text_color=WHITE, dot_color=GOLD,
                     width=None, height=None, text_size=14,
                     shadow=False):
    """Small 'POLL' pill, right-aligned.

    position: 'top-right' (default, just below the top bar) or
              'bottom-right' (bottom band aligned with discussion break).
    width / height: override the default 1.05" × 0.42" pill size.
    text_size: "POLL" font size (pt).  Scales up when the pill is enlarged.
    dot_color: pass None to skip the leading accent dot.
    shadow: add a soft drop shadow to the pill (matches discussion-break
        chrome).  Default False to preserve the original flat top-right
        pill look.
    """
    pill_w = width if width is not None else Inches(1.05)
    pill_h = height if height is not None else Inches(0.42)
    pill_x = SLIDE_W - MARGIN - pill_w
    if position == 'bottom-right':
        # Bottom-aligned with where _add_discussion_break ends
        # (top=6.25, height=0.72 → bottom=6.97).  Pill height varies, so
        # pick top so the pill BOTTOM lands on the same 6.97 baseline.
        pill_y = Inches(6.97) - pill_h
    else:
        pill_y = Inches(0.55)

    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  pill_x, pill_y, pill_w, pill_h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    if shadow:
        _add_drop_shadow(shp)
    else:
        shp.shadow.inherit = False
    try:
        shp.adjustments[0] = 0.5
    except Exception:
        pass
    _add_text(slide, pill_x, pill_y, pill_w, pill_h,
              "POLL", size=text_size, bold=True, color=text_color, font="Calibri",
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if dot_color is not None:
        dot_size = Inches(0.16)
        dot_x = pill_x - Inches(0.16) - dot_size
        dot_y = pill_y + (pill_h - dot_size) // 2
        _add_rect(slide, dot_x, dot_y, dot_size, dot_size, dot_color)


def make_poll_slide(prs, page_num, section_tag, title, options, *,
                    instructions="Respond at PollEv.com/nvoigtlaender",
                    size=30, line_spacing_pts=22):
    """Layout 5 – Poll slide.

    options: list of strings (A./B./C./D. auto-numbering applied).
    """
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, section_tag)
    _draw_action_title(slide, title)
    _draw_poll_pill(slide)

    _add_bulleted_list(
        slide,
        left=MARGIN,
        top=Inches(2.0),
        width=RULE_W,
        height=Inches(4.4),
        items=options,
        size=size, color=NAVY, bullet_color=NAVY,
        line_spacing_pts=line_spacing_pts,
        autonum_scheme='alphaUcPeriod',
    )

    _add_text(slide, MARGIN, Inches(6.55), RULE_W, Inches(0.4),
              instructions, size=16, italic=True, color=GRAY,
              align=PP_ALIGN.RIGHT)

    _draw_footer(slide, FOOTER_TEXT, page_num)
    return slide


# --------------------------------------------------------------------------
# Slide-jump hyperlinks — make a run or shape clickable in slideshow
# mode so it advances PowerPoint to a different slide.  python-pptx
# doesn't expose this directly, so we manipulate the rPr / cNvPr XML
# and register the slide relationship via the public OPC API.
# --------------------------------------------------------------------------

def _add_slide_jump_hyperlink_run(source_slide, run, target_slide,
                                    *, lock_color=True, underline=True):
    """Make ``run`` a clickable hyperlink that advances to ``target_slide``.

    ``lock_color=True`` adds the Office hyperlinkcolor extension so the
    run's explicit ``solidFill`` color is preserved at render time —
    without it, PowerPoint repaints hyperlink text in the theme's
    ``hlink`` color (a light blue) regardless of what the rPr says.

    ``underline=True`` sets ``u="sng"`` so the hyperlink stays visibly
    underlined even when the color-lock suppresses the theme styling.
    """
    rId = source_slide.part.relate_to(target_slide.part, RT.SLIDE)
    rPr = run._r.get_or_add_rPr()
    if underline:
        rPr.set('u', 'sng')
    for hl in rPr.findall(qn('a:hlinkClick')):
        rPr.remove(hl)
    hlinkClick = ET.SubElement(rPr, qn('a:hlinkClick'))
    hlinkClick.set(qn('r:id'), rId)
    hlinkClick.set('action', 'ppaction://hlinksldjump')
    if lock_color:
        AHYP_NS = 'http://schemas.microsoft.com/office/drawing/2018/hyperlinkcolor'
        ext_xml = (
            f'<a:extLst xmlns:a="{A_NS}">'
            f'<a:ext uri="{{A12FA001-AC4F-418D-AE19-62706E023703}}">'
            f'<ahyp:hlinkClr xmlns:ahyp="{AHYP_NS}" val="tx"/>'
            f'</a:ext>'
            f'</a:extLst>'
        )
        hlinkClick.append(ET.fromstring(ext_xml))


def _add_slide_jump_hyperlink_shape(source_slide, shape, target_slide):
    """Make ``shape`` clickable so the whole shape jumps to ``target_slide``."""
    rId = source_slide.part.relate_to(target_slide.part, RT.SLIDE)
    nvSpPr = shape._element.find(qn('p:nvSpPr'))
    cNvPr = nvSpPr.find(qn('p:cNvPr')) if nvSpPr is not None else None
    if cNvPr is None:
        return
    for hl in cNvPr.findall(qn('a:hlinkClick')):
        cNvPr.remove(hl)
    hlinkClick = ET.SubElement(cNvPr, qn('a:hlinkClick'))
    hlinkClick.set(qn('r:id'), rId)
    hlinkClick.set('action', 'ppaction://hlinksldjump')


# --------------------------------------------------------------------------
# Image helpers — embed source-deck images into the new deck.
# Images are pre-extracted to _source_images/slide{N}_{rId}.{ext}.
# --------------------------------------------------------------------------

SRC_IMG_DIR = Path(__file__).parent / "_source_images"


def _add_drop_shadow(shape, *, blur="50800", dist="38100",
                      direction="2700000", alpha="45000"):
    """Add a soft drop shadow to any shape that exposes spPr (pictures,
    rectangles, rounded rects).  Default: 4pt blur, 3pt offset, 45° down-
    right, 45% opacity black.  Used deck-wide for figures/boxes."""
    try:
        spPr = shape._element.spPr
    except AttributeError:
        # Fallback for shapes whose XML element exposes spPr only via find()
        spPr = shape._element.find(qn('p:spPr'))
    if spPr is None:
        return shape
    for old in spPr.findall(qn('a:effectLst')):
        spPr.remove(old)
    effLst = ET.SubElement(spPr, qn('a:effectLst'))
    outerShdw = ET.SubElement(effLst, qn('a:outerShdw'))
    outerShdw.set('blurRad', str(blur))
    outerShdw.set('dist', str(dist))
    outerShdw.set('dir', str(direction))
    outerShdw.set('algn', 'tl')
    outerShdw.set('rotWithShape', '0')
    rgb = ET.SubElement(outerShdw, qn('a:srgbClr'))
    rgb.set('val', '000000')
    a = ET.SubElement(rgb, qn('a:alpha'))
    a.set('val', str(alpha))
    return shape


def _add_graphicframe_shadow(slide, left, top, width, height, *,
                              shadow_alpha=45000):
    """White backing rectangle with an outerShdw effect, behind a
    graphicFrame (table or chart).  graphicFrames can't host
    a:effectLst directly; this rect supplies the shadow projected
    OUTSIDE its bounds.

    Charts have transparent plot areas by default, so a coloured backing
    bleeds through and tints the figure.  Use white instead — the chart
    sees white through its own transparent areas (clean background) and
    the shadow renders only at the visible edges.  Call BEFORE adding
    the table/chart so z-order is correct.
    """
    shdw = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        int(left), int(top), int(width), int(height),
    )
    shdw.fill.solid()
    shdw.fill.fore_color.rgb = WHITE
    shdw.line.fill.background()
    shdw.shadow.inherit = False
    sp_pr = shdw._element.spPr
    # Strip any default <a:effectLst/> python-pptx may have inserted
    # (duplicate effectLst makes PowerPoint refuse to open the file).
    for old in sp_pr.findall(qn('a:effectLst')):
        sp_pr.remove(old)
    effLst = ET.SubElement(sp_pr, qn('a:effectLst'))
    outerShdw = ET.SubElement(effLst, qn('a:outerShdw'))
    outerShdw.set('blurRad', '50800')
    outerShdw.set('dist', '38100')
    outerShdw.set('dir', '2700000')
    outerShdw.set('algn', 'tl')
    outerShdw.set('rotWithShape', '0')
    rgb = ET.SubElement(outerShdw, qn('a:srgbClr'))
    rgb.set('val', '000000')
    a = ET.SubElement(rgb, qn('a:alpha'))
    a.set('val', str(int(shadow_alpha)))
    return shdw


def _add_source_image(slide, src_slide_no, rid, *, left, top, width=None,
                      height=None, shadow=True):
    """Place a source-deck image on the new slide.

    `shadow=True` (default) adds a soft drop shadow so figures pop off the
    background — applied deck-wide per the latest visual direction.  Set
    `shadow=False` for niche cases (transparent PNGs, screenshots that
    already include a shadow, etc.).
    """
    candidates = list(SRC_IMG_DIR.glob(f"slide{src_slide_no}_{rid}.*"))
    if not candidates:
        return None
    img = candidates[0]
    kwargs = {"left": int(left), "top": int(top)}
    if width is not None:
        kwargs["width"] = int(width)
    if height is not None:
        kwargs["height"] = int(height)
    pic = slide.shapes.add_picture(str(img), **kwargs)
    if shadow:
        _add_drop_shadow(pic)
    return pic


def _apply_picture_style(pic, *, corner_pct=8,
                          shadow_blur=50800, shadow_dist=38100,
                          shadow_dir=2700000, shadow_alpha=50000):
    """Apply rounded corners + drop shadow to a picture shape.

    corner_pct: rounded-corner radius as percent of shorter side (8 ≈ subtle).
    shadow_blur/dist: EMU; defaults give a soft 4pt blur, 3pt offset.
    shadow_dir: 2700000 = 45° down-right (standard).
    shadow_alpha: 50000 = 50% opacity black shadow.
    """
    spPr = pic._element.find(qn('p:spPr'))
    if spPr is None:
        return pic
    # Replace any existing prstGeom with roundRect at corner_pct
    for old in spPr.findall(qn('a:prstGeom')):
        spPr.remove(old)
    prstGeom = ET.Element(qn('a:prstGeom'))
    prstGeom.set('prst', 'roundRect')
    avLst = ET.SubElement(prstGeom, qn('a:avLst'))
    gd = ET.SubElement(avLst, qn('a:gd'))
    gd.set('name', 'adj')
    gd.set('fmla', f'val {int(corner_pct * 1000)}')
    # prstGeom must come after a:xfrm
    xfrm = spPr.find(qn('a:xfrm'))
    if xfrm is not None:
        xfrm.addnext(prstGeom)
    else:
        spPr.insert(0, prstGeom)
    # Replace any existing effectLst with a single outer shadow
    for old in spPr.findall(qn('a:effectLst')):
        spPr.remove(old)
    effectLst = ET.SubElement(spPr, qn('a:effectLst'))
    outerShdw = ET.SubElement(effectLst, qn('a:outerShdw'))
    outerShdw.set('blurRad', str(int(shadow_blur)))
    outerShdw.set('dist', str(int(shadow_dist)))
    outerShdw.set('dir', str(int(shadow_dir)))
    outerShdw.set('algn', 'tl')
    outerShdw.set('rotWithShape', '0')
    rgb = ET.SubElement(outerShdw, qn('a:srgbClr'))
    rgb.set('val', '000000')
    alpha = ET.SubElement(rgb, qn('a:alpha'))
    alpha.set('val', str(int(shadow_alpha)))
    return pic


# --------------------------------------------------------------------------
# Illustrative callout boxes — the "major-concept" badges, "Teaching Note"
# bars, and bottom-of-slide takeaway bands that recur throughout the source
# deck.  Replicating them faithfully (in template colors) is what makes the
# slides feel like the original, not a stripped-down rewrite.
# --------------------------------------------------------------------------

def _add_takeaway_bar(slide, text, *, top=Inches(6.4), width=None,
                       height=Inches(0.55), left=None,
                       fill=GOLD, text_color=WHITE,
                       size=20, font="Calibri", bold=True,
                       rounded=False, shadow=False):
    """Bottom-of-slide takeaway band — the 'major concept' callout.

    rounded=True  → ROUNDED_RECTANGLE with ~30% corner radius.
    shadow=True   → soft drop shadow (off by default to keep the flat
                    look on the majority of slides).
    left=None     → horizontally centered. Pass an EMU/Inches value to
                    pin the bar's left edge (used for hand-positioned
                    takeaways).
    """
    if width is None:
        width = Inches(9.6)
    if left is None:
        left = (SLIDE_W - width) // 2
    if not (rounded or shadow):
        return _add_filled_box(slide, left, top, width, height, text,
                                fill=fill, text_color=text_color,
                                size=size, bold=bold, font=font)
    # Inline build for the rounded / shadowed variant
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, int(left), int(top), int(width), int(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    if rounded:
        try: shape.adjustments[0] = 0.30
        except Exception: pass
    shape.shadow.inherit = False
    if shadow:
        _add_drop_shadow(shape)
    tf = shape.text_frame
    tf.margin_left = tf.margin_right = Inches(0.1)
    tf.margin_top = tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = text_color
    return shape


def _add_teaching_note(slide, text, *, top=Inches(6.6), width=None,
                        height=Inches(0.6), left=None,
                        rounded=False, pdf_icon=False, label_color=GOLD,
                        fill_rgb=None):
    """External-document reference card.

    Visually distinct from in-slide callouts: cream/parchment fill, navy
    DASHED border, italic navy text, with a small page-icon glyph on the
    left — clearly signals 'this links to an external Teaching Note doc'.

    Optional styling for the slide-32 (page 33) treatment:
      • ``left``       — pin to an absolute X (default: horizontally center)
      • ``rounded``    — rounded corners + soft drop shadow (default flat)
      • ``pdf_icon``   — render the page glyph as a white folded-corner
                        with bold "PDF" text, sized to nearly fill the card
                        (default: small navy folded-corner, no text)
      • ``label_color``— color of the "SEE TEACHING NOTE →" prefix
                        (default GOLD; pass NAVY for the page-33 look).
      • ``fill_rgb``   — override the cream/parchment card fill
                        (default RGBColor(0xF4, 0xF1, 0xEA)).
    """
    if width is None:
        width = Inches(8.0)
    if left is None:
        left = int((SLIDE_W - width) // 2)
    left, top, width, height = int(left), int(top), int(width), int(height)

    # Card: cream fill, dashed navy border.  ``rounded=True`` switches the
    # base shape to ROUNDED_RECTANGLE and adds a soft drop shadow.
    base_shape = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    card = slide.shapes.add_shape(base_shape, left, top, width, height)
    if rounded:
        try: card.adjustments[0] = 0.20
        except Exception: pass
    card.fill.solid()
    card.fill.fore_color.rgb = (fill_rgb if fill_rgb is not None
                                  else RGBColor(0xF4, 0xF1, 0xEA))  # parchment cream default
    card.line.color.rgb = NAVY
    card.line.width = Pt(1.25)
    card.shadow.inherit = False
    if rounded:
        _add_drop_shadow(card)
    # Dashed line style
    ln = card.line._get_or_add_ln()
    # Remove any existing prstDash
    for el in ln.findall(qn('a:prstDash')):
        ln.remove(el)
    prstDash = ET.SubElement(ln, qn('a:prstDash'))
    prstDash.set('val', 'dash')

    # Empty text on the card; we add the icon + text as separate textboxes
    tf = card.text_frame
    tf.text = ""

    # Small "page" icon on the left – a folded-corner shape.  In default
    # mode it's a small filled-navy glyph (logo-ish); in pdf_icon mode it's
    # a larger WHITE sheet with thin navy border + bold "PDF" text so it
    # reads as a generic PDF document at a glance.
    icon_size = Inches(0.5) if pdf_icon else Inches(0.4)
    icon_x = left + Inches(0.2)
    icon_y = top + (height - icon_size) // 2
    page = slide.shapes.add_shape(MSO_SHAPE.FOLDED_CORNER,
                                   int(icon_x), int(icon_y),
                                   int(icon_size), int(icon_size))
    page.fill.solid()
    if pdf_icon:
        page.fill.fore_color.rgb = WHITE
        page.line.color.rgb = NAVY
        page.line.width = Pt(0.75)
    else:
        page.fill.fore_color.rgb = NAVY
        page.line.fill.background()
    page.shadow.inherit = False
    if pdf_icon:
        # "PDF" label inside the white sheet.
        ptf = page.text_frame
        ptf.word_wrap = False
        ptf.margin_left = 0
        ptf.margin_right = 0
        ptf.margin_top = 0
        ptf.margin_bottom = 0
        ptf.vertical_anchor = MSO_ANCHOR.MIDDLE
        pp = ptf.paragraphs[0]
        pp.alignment = PP_ALIGN.CENTER
        pr = pp.add_run()
        pr.text = "PDF"
        pr.font.name = "Calibri"
        pr.font.size = Pt(11)
        pr.font.bold = True
        pr.font.color.rgb = NAVY

    # Label text (italic, navy) inside the card to the right of the icon
    txt_left = int(icon_x + icon_size + Inches(0.2))
    txt_w = int(width - (icon_x + icon_size + Inches(0.4) - left))
    label_box = slide.shapes.add_textbox(txt_left, top,
                                          txt_w, height)
    label_tf = label_box.text_frame
    label_tf.word_wrap = True
    label_tf.margin_left = 0
    label_tf.margin_right = 0
    label_tf.margin_top = 0
    label_tf.margin_bottom = 0
    label_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = label_tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    # First run: "See teaching note:" prefix in caller-supplied color
    # (GOLD by default; NAVY for the slide-33 treatment).
    r1 = p.add_run()
    r1.text = "SEE TEACHING NOTE  →  "
    r1.font.name = "Calibri"
    r1.font.size = Pt(12)
    r1.font.bold = True
    r1.font.color.rgb = label_color
    # Second run: the title of the note, italic navy
    r2 = p.add_run()
    r2.text = text
    r2.font.name = "Calibri"
    r2.font.size = Pt(16)
    r2.font.italic = True
    r2.font.bold = True
    r2.font.color.rgb = NAVY
    return card


def _add_discussion_break(slide, *, top=Inches(6.25), width=Inches(4.8),
                           left=None, text="Discussion Break"):
    """Rounded-parallelogram 'discussion break' badge (bottom-right).

    Custom-geometry shape: top and bottom edges are horizontal; the left
    and right edges slant at 45° in real space (skew = height of the
    shape).  All four corners are slightly rounded.  Gold fill, navy
    bold text, soft drop shadow.

    left=None → pin to the right edge with the default MARGIN. Pass an
    EMU/Inches value to override (used for hand-positioned badges).
    """
    height = Inches(0.72)
    if left is None:
        left = SLIDE_W - MARGIN - width
    left, top, width, height = int(left), int(top), int(width), int(height)
    # Compute skew in path-coordinate units so that left/right sides slant
    # at 45° in REAL space:  horizontal offset of top edge == shape height.
    skew = int(100000 * height / width) if width else 15000
    skew = min(max(skew, 6000), 35000)
    r = 5000          # corner radius in path units — gentle rounding

    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = GOLD
    shp.line.fill.background()
    shp.shadow.inherit = False
    # Strip the default prstGeom – we'll inject custGeom instead.
    spPr = shp._element.spPr
    for old in spPr.findall(qn('a:prstGeom')):
        spPr.remove(old)
    rs = int(r * skew / 100000)
    # IMPORTANT: <a:rect> defines the TEXT bounding rectangle inside the
    # custom geometry.  Set it to the parallelogram's inscribed rectangle
    # (from TL vertex to BR vertex) so PowerPoint won't render text past
    # the slanted edges, regardless of the text frame's own lIns/rIns.
    custgeom_xml = (
        f'<a:custGeom xmlns:a="{A_NS}">'
        f'<a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/>'
        f'<a:rect l="{skew}" t="0" r="{100000-skew}" b="100000"/>'
        f'<a:pathLst><a:path w="100000" h="100000">'
        # Top-left rounded corner (vertex at (skew, 0))
        f'<a:moveTo><a:pt x="{skew+r}" y="0"/></a:moveTo>'
        # Top edge
        f'<a:lnTo><a:pt x="{100000-r}" y="0"/></a:lnTo>'
        # Top-right corner round
        f'<a:cubicBezTo>'
        f'<a:pt x="100000" y="0"/><a:pt x="100000" y="0"/>'
        f'<a:pt x="{100000-rs}" y="{r}"/>'
        f'</a:cubicBezTo>'
        # Right slanted side (down-left)
        f'<a:lnTo><a:pt x="{100000-skew+rs}" y="{100000-r}"/></a:lnTo>'
        # Bottom-right corner round (vertex at (100000-skew, 100000))
        f'<a:cubicBezTo>'
        f'<a:pt x="{100000-skew}" y="100000"/>'
        f'<a:pt x="{100000-skew}" y="100000"/>'
        f'<a:pt x="{100000-skew-r}" y="100000"/>'
        f'</a:cubicBezTo>'
        # Bottom edge
        f'<a:lnTo><a:pt x="{r}" y="100000"/></a:lnTo>'
        # Bottom-left corner round (vertex at (0, 100000))
        f'<a:cubicBezTo>'
        f'<a:pt x="0" y="100000"/><a:pt x="0" y="100000"/>'
        f'<a:pt x="{rs}" y="{100000-r}"/>'
        f'</a:cubicBezTo>'
        # Left slanted side (up-right)
        f'<a:lnTo><a:pt x="{skew-rs}" y="{r}"/></a:lnTo>'
        # Top-left corner round (close the path)
        f'<a:cubicBezTo>'
        f'<a:pt x="{skew}" y="0"/><a:pt x="{skew}" y="0"/>'
        f'<a:pt x="{skew+r}" y="0"/>'
        f'</a:cubicBezTo>'
        f'<a:close/>'
        f'</a:path></a:pathLst>'
        f'</a:custGeom>'
    )
    custgeom = ET.fromstring(custgeom_xml)
    # Insert custGeom right after a:xfrm (schema order)
    xfrm = spPr.find(qn('a:xfrm'))
    if xfrm is not None:
        xfrm.addnext(custgeom)
    else:
        spPr.insert(0, custgeom)

    # Drop shadow (45° down-right, 50% opacity).
    for old in spPr.findall(qn('a:effectLst')):
        spPr.remove(old)
    effectLst = ET.SubElement(spPr, qn('a:effectLst'))
    outerShdw = ET.SubElement(effectLst, qn('a:outerShdw'))
    outerShdw.set('blurRad', '50800')
    outerShdw.set('dist', '38100')
    outerShdw.set('dir', '2700000')
    outerShdw.set('algn', 'tl')
    outerShdw.set('rotWithShape', '0')
    rgb = ET.SubElement(outerShdw, qn('a:srgbClr'))
    rgb.set('val', '000000')
    alpha = ET.SubElement(rgb, qn('a:alpha'))
    alpha.set('val', '50000')

    # Text — render as a SEPARATE textbox overlaid on top of the
    # parallelogram, positioned exactly inside the inscribed rectangle.
    # This decouples text placement from the shape geometry and avoids
    # PowerPoint rendering the run past the slanted edges (which can
    # happen with the in-shape text frame on some PowerPoint versions
    # even with <a:rect> set inside custGeom).
    # In real EMU, the skew equals the shape height (45° slant), so the
    # inscribed rectangle spans (left + height) → (left + width - height).
    skew_emu = height
    ins_left = left + skew_emu
    ins_top = top
    ins_w = width - 2 * skew_emu
    ins_h = height
    txt = slide.shapes.add_textbox(int(ins_left), int(ins_top),
                                     int(ins_w), int(ins_h))
    tf = txt.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(28)         # bumped 20 → 28 on 2026-05-16
    run.font.bold = True
    run.font.color.rgb = NAVY
    return shp


def _add_callout_box(slide, left, top, width, height, text, *,
                      fill=GOLD, text_color=WHITE, size=14, bold=True):
    """Small free-form annotation/callout (e.g., 'plot the slope', 'Revenue
    per car net of material cost').  Used to mark a graph or sub-region."""
    return _add_filled_box(slide, left, top, width, height, text,
                            fill=fill, text_color=text_color,
                            size=size, bold=bold, font="Calibri")


def _add_anchor_burst(slide, left, top, width, height,
                       top_text, bottom_text=None, extra_text=None,
                       *, fill=GOLD, text_color=NAVY,
                       top_size=14, bottom_size=10):
    """12-point star background + a separate text-box overlay.

    The star is purely decorative; the text lives in a normal
    rectangular text box layered on top so the text isn't constrained
    by the star's geometry (no risk of bleeding into the points).
    Reusable on every slide where MB = MC is being invoked, so the same
    visual pattern carries over.
    """
    left, top, width, height = int(left), int(top), int(width), int(height)

    # 1. Decorative star background (no text)
    shp = slide.shapes.add_shape(MSO_SHAPE.STAR_12_POINT, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = NAVY
    shp.line.width = Pt(1.0)
    shp.shadow.inherit = False
    # Soft drop shadow — added 2026-05-15 per user request, so the MB=MC
    # star reads as lifted off the slide like every other content shape.
    _add_drop_shadow(shp)
    # Suppress any auto-inserted text frame contents on the shape itself.
    shp.text_frame.text = ""

    # 2. Overlay text box, sized to the star's inscribed body
    inner_w = int(width * 0.65)
    inner_h = int(height * 0.55)
    inner_x = left + (width - inner_w) // 2
    inner_y = top + (height - inner_h) // 2

    box = slide.shapes.add_textbox(inner_x, inner_y, inner_w, inner_h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = top_text
    r1.font.name = 'Calibri'
    r1.font.size = Pt(top_size)
    r1.font.bold = True
    r1.font.color.rgb = text_color
    if bottom_text:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = bottom_text
        r2.font.name = 'Calibri'
        r2.font.size = Pt(bottom_size)
        r2.font.italic = True
        r2.font.bold = True
        r2.font.color.rgb = text_color
    if extra_text:
        p3 = tf.add_paragraph()
        p3.alignment = PP_ALIGN.CENTER
        r3 = p3.add_run()
        r3.text = extra_text
        r3.font.name = 'Calibri'
        r3.font.size = Pt(bottom_size)
        r3.font.italic = True
        r3.font.bold = True
        r3.font.color.rgb = text_color
    return shp


# --------------------------------------------------------------------------
# OMML (Office Math Markup Language) equation helper – gives formulas a
# proper TeX-style render with italic variables, stacked fractions, real
# subscripts/superscripts.  Uses Cambria Math (the standard PPT math font).
# --------------------------------------------------------------------------

M_NS = 'http://schemas.openxmlformats.org/officeDocument/2006/math'
A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
A14_NS = 'http://schemas.microsoft.com/office/drawing/2010/main'


def _omml_fill(color):
    """Build the inner ``<a:solidFill>`` clause for an OMML run's rPr.

    ``color`` may be an ``RGBColor`` instance (or a 3-tuple of ints).
    Returns an empty string when ``color`` is None so callers can splice
    the result into rPr unconditionally.
    """
    if color is None:
        return ''
    return (
        f'<a:solidFill><a:srgbClr val="'
        f'{color[0]:02X}{color[1]:02X}{color[2]:02X}'
        f'"/></a:solidFill>'
    )


def _omml_run(text, *, color=None):
    """OMML run for an italic variable (default math style).

    Inside an oMath, italic style is the math default for Latin letters;
    we leave m:rPr out entirely so the Cambria Math italic comes through.
    The a:rPr applies drawing-level font sizing/coloring.  Pass ``color``
    to tint the run (e.g., green ΔL / ΔQ in the slide-14 Convention box).
    """
    return (
        f'<m:r xmlns:m="{M_NS}">'
        f'<a:rPr xmlns:a="{A_NS}" lang="en-US" b="0" i="1">'
        f'{_omml_fill(color)}'
        f'<a:latin typeface="Cambria Math"/>'
        f'<a:ea typeface="Cambria Math"/>'
        f'</a:rPr>'
        f'<m:t>{text}</m:t>'
        f'</m:r>'
    )


def _omml_text(text, *, color=None):
    """Upright-style OMML run (for operators, numbers, acronyms).

    Force plain (upright) style via <m:rPr><m:sty m:val="p"/></m:rPr> – this
    is the documented way to disable the math-default italics for the
    enclosed run.  Pass ``color`` to tint the run.
    """
    return (
        f'<m:r xmlns:m="{M_NS}">'
        f'<m:rPr><m:sty m:val="p"/></m:rPr>'
        f'<a:rPr xmlns:a="{A_NS}" lang="en-US" b="0" i="0">'
        f'{_omml_fill(color)}'
        f'<a:latin typeface="Cambria Math"/>'
        f'</a:rPr>'
        f'<m:t xml:space="preserve">{text}</m:t>'
        f'</m:r>'
    )


def _omml_sub(base, sub):
    """OMML subscript: base with subscript expression."""
    return (
        f'<m:sSub xmlns:m="{M_NS}">'
        f'<m:sSubPr><m:ctrlPr>'
        f'<a:rPr xmlns:a="{A_NS}" lang="en-US" i="1">'
        f'<a:latin typeface="Cambria Math"/></a:rPr>'
        f'</m:ctrlPr></m:sSubPr>'
        f'<m:e>{base}</m:e>'
        f'<m:sub>{sub}</m:sub>'
        f'</m:sSub>'
    )


def _omml_frac(num, den):
    """OMML stacked fraction: num / den."""
    return (
        f'<m:f xmlns:m="{M_NS}">'
        f'<m:fPr><m:ctrlPr>'
        f'<a:rPr xmlns:a="{A_NS}" lang="en-US" i="1">'
        f'<a:latin typeface="Cambria Math"/></a:rPr>'
        f'</m:ctrlPr></m:fPr>'
        f'<m:num>{num}</m:num>'
        f'<m:den>{den}</m:den>'
        f'</m:f>'
    )


def _omml_sup(base, sup):
    """OMML superscript: base^sup (e.g. Q²)."""
    return (
        f'<m:sSup xmlns:m="{M_NS}">'
        f'<m:sSupPr><m:ctrlPr>'
        f'<a:rPr xmlns:a="{A_NS}" lang="en-US" i="1">'
        f'<a:latin typeface="Cambria Math"/></a:rPr>'
        f'</m:ctrlPr></m:sSupPr>'
        f'<m:e>{base}</m:e>'
        f'<m:sup>{sup}</m:sup>'
        f'</m:sSup>'
    )


def _add_dashed_gridlines(axis_el):
    """Dashed light-grey major gridlines on an axis (Cobb-Douglas style)."""
    for old in axis_el.findall(qn('c:majorGridlines')):
        axis_el.remove(old)
    gl = ET.Element(qn('c:majorGridlines'))
    sp = ET.SubElement(gl, qn('c:spPr'))
    ln = ET.SubElement(sp, qn('a:ln'))
    ln.set('w', '9525')
    ln.set('cap', 'flat'); ln.set('cmpd', 'sng'); ln.set('algn', 'ctr')
    fill = ET.SubElement(ln, qn('a:solidFill'))
    clr = ET.SubElement(fill, qn('a:srgbClr')); clr.set('val', 'C8CDD3')
    dash = ET.SubElement(ln, qn('a:prstDash')); dash.set('val', 'dash')
    axpos = axis_el.find(qn('c:axPos'))
    if axpos is not None:
        axpos.addnext(gl)


def _align_x_labels_with_ticks(value_axis):
    """Set <c:crossBetween val="midCat"/> on a value axis so the category
    labels and tick marks align (default for line charts in OOXML is
    "between", which leaves labels visually between adjacent ticks).
    """
    val_el = value_axis._element
    for old in val_el.findall(qn('c:crossBetween')):
        val_el.remove(old)
    cb = ET.Element(qn('c:crossBetween'))
    cb.set('val', 'midCat')
    # Schema position: c:crossBetween follows c:crosses(At) and precedes
    # c:majorUnit.  Insert before c:majorUnit if present; else append.
    mu_el = val_el.find(qn('c:majorUnit'))
    if mu_el is not None:
        mu_el.addprevious(cb)
    else:
        val_el.append(cb)


def _make_simple_line_chart(slide, x, y, w, h, categories, values, *,
                              line_color, x_title, y_title,
                              y_min=0, y_max=None, y_unit=None,
                              marker='circle'):
    """Single-series line+markers chart with dashed light-grey gridlines.

    Same visual conventions as slide 11 (Calibri navy labels, dashed C8CDD3
    major gridlines, marker size 7) but no legend/title.
    """
    cd = CategoryChartData()
    cd.categories = list(categories)
    cd.add_series("Y", values)
    # Drop-shadow rectangle behind the chart (graphicFrames can't host shadow).
    _add_graphicframe_shadow(slide, x, y, w, h)
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE,
        int(x), int(y), int(w), int(h), cd,
    )
    chart = chart_shape.chart
    chart.has_title = False
    chart.has_legend = False

    clr_hex = f'{line_color[0]:02X}{line_color[1]:02X}{line_color[2]:02X}'

    for series in chart.series:
        line = series.format.line
        line.color.rgb = line_color
        line.width = Pt(2.5)
        ser_xml = series._element
        for old in ser_xml.findall(qn('c:marker')):
            ser_xml.remove(old)
        m = ET.SubElement(ser_xml, qn('c:marker'))
        sym = ET.SubElement(m, qn('c:symbol')); sym.set('val', marker)
        sz_el = ET.SubElement(m, qn('c:size')); sz_el.set('val', '7')
        sp = ET.SubElement(m, qn('c:spPr'))
        fl = ET.SubElement(sp, qn('a:solidFill'))
        rg = ET.SubElement(fl, qn('a:srgbClr')); rg.set('val', clr_hex)
        ln = ET.SubElement(sp, qn('a:ln'))
        lf = ET.SubElement(ln, qn('a:solidFill'))
        lr = ET.SubElement(lf, qn('a:srgbClr')); lr.set('val', clr_hex)
        # disable smoothing (straight segments between points)
        for sm in ser_xml.findall(qn('c:smooth')):
            ser_xml.remove(sm)
        smooth = ET.SubElement(ser_xml, qn('c:smooth'))
        smooth.set('val', '0')

    # Axes – axis titles in BOLD ITALIC navy (per course CLAUDE.md);
    # tick labels in regular Calibri navy.
    cat = chart.category_axis
    cat.tick_labels.font.name = "Calibri"
    cat.tick_labels.font.size = Pt(10)
    cat.tick_labels.font.color.rgb = NAVY
    cat.has_title = True
    cat.axis_title.text_frame.text = x_title
    ar = cat.axis_title.text_frame.paragraphs[0].runs[0]
    ar.font.name = "Calibri"; ar.font.size = Pt(12)
    ar.font.bold = True; ar.font.italic = True
    ar.font.color.rgb = NAVY

    val = chart.value_axis
    val.tick_labels.font.name = "Calibri"
    val.tick_labels.font.size = Pt(10)
    val.tick_labels.font.color.rgb = NAVY
    val.has_title = True
    val.axis_title.text_frame.text = y_title
    ar = val.axis_title.text_frame.paragraphs[0].runs[0]
    ar.font.name = "Calibri"; ar.font.size = Pt(12)
    ar.font.bold = True; ar.font.italic = True
    ar.font.color.rgb = NAVY
    if y_max is not None:
        val.minimum_scale = y_min
        val.maximum_scale = y_max
    if y_unit is not None:
        val.major_unit = y_unit

    # Align X-axis category labels with tick marks (default OOXML places
    # them in the gaps between ticks).
    _align_x_labels_with_ticks(val)

    _add_dashed_gridlines(cat._element)
    _add_dashed_gridlines(val._element)
    return chart_shape


def _make_multi_line_chart(slide, x, y, w, h, categories, series, *,
                             x_title, y_title,
                             y_min=0, y_max=None, y_unit=None,
                             legend=True, legend_pos=('0.08', '0.10', '0.20', '0.20')):
    """Multi-series line+markers chart with the deck's standard styling.

    series: list of (name, values, color: RGBColor, marker: str) tuples.
    legend_pos: (x, y, w, h) in chart-fraction units (str) – top-left default.
    """
    _add_graphicframe_shadow(slide, x, y, w, h)
    cd = CategoryChartData()
    cd.categories = list(categories)
    for name, values, _color, _marker in series:
        cd.add_series(name, values)
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE,
        int(x), int(y), int(w), int(h), cd,
    )
    chart = chart_shape.chart
    chart.has_title = False
    chart.has_legend = legend

    # Per-series styling: line color + marker
    for idx, ser in enumerate(chart.series):
        name, values, color, marker = series[idx]
        line = ser.format.line
        line.color.rgb = color
        line.width = Pt(2.5)
        ser_xml = ser._element
        clr_hex = f'{color[0]:02X}{color[1]:02X}{color[2]:02X}'
        for old in ser_xml.findall(qn('c:marker')):
            ser_xml.remove(old)
        m = ET.SubElement(ser_xml, qn('c:marker'))
        sym = ET.SubElement(m, qn('c:symbol')); sym.set('val', marker)
        sz_el = ET.SubElement(m, qn('c:size')); sz_el.set('val', '7')
        sp = ET.SubElement(m, qn('c:spPr'))
        fl = ET.SubElement(sp, qn('a:solidFill'))
        rg = ET.SubElement(fl, qn('a:srgbClr')); rg.set('val', clr_hex)
        ln = ET.SubElement(sp, qn('a:ln'))
        lf = ET.SubElement(ln, qn('a:solidFill'))
        lr = ET.SubElement(lf, qn('a:srgbClr')); lr.set('val', clr_hex)
        for sm in ser_xml.findall(qn('c:smooth')):
            ser_xml.remove(sm)
        smooth = ET.SubElement(ser_xml, qn('c:smooth'))
        smooth.set('val', '0')

    # Axes – bold italic navy titles per course style
    cat = chart.category_axis
    cat.tick_labels.font.name = "Calibri"
    cat.tick_labels.font.size = Pt(10)
    cat.tick_labels.font.color.rgb = NAVY
    cat.has_title = True
    cat.axis_title.text_frame.text = x_title
    ar = cat.axis_title.text_frame.paragraphs[0].runs[0]
    ar.font.name = "Calibri"; ar.font.size = Pt(12)
    ar.font.bold = True; ar.font.italic = True; ar.font.color.rgb = NAVY

    val = chart.value_axis
    val.tick_labels.font.name = "Calibri"
    val.tick_labels.font.size = Pt(10)
    val.tick_labels.font.color.rgb = NAVY
    val.has_title = True
    val.axis_title.text_frame.text = y_title
    ar = val.axis_title.text_frame.paragraphs[0].runs[0]
    ar.font.name = "Calibri"; ar.font.size = Pt(12)
    ar.font.bold = True; ar.font.italic = True; ar.font.color.rgb = NAVY
    if y_max is not None:
        val.minimum_scale = y_min
        val.maximum_scale = y_max
    if y_unit is not None:
        val.major_unit = y_unit

    # Align X-axis category labels with tick marks.
    _align_x_labels_with_ticks(val)

    _add_dashed_gridlines(cat._element)
    _add_dashed_gridlines(val._element)

    # Legend top-left inside plot with white fill
    if legend:
        leg_el = chart.legend._element
        chart.legend.font.name = "Calibri"
        chart.legend.font.size = Pt(11)
        chart.legend.font.color.rgb = NAVY
        chart.legend.include_in_layout = False
        # Strip default legendPos / layout, replace
        for old in leg_el.findall(qn('c:layout')):
            leg_el.remove(old)
        for old in leg_el.findall(qn('c:legendPos')):
            leg_el.remove(old)
        pos = ET.SubElement(leg_el, qn('c:legendPos')); pos.set('val', 'tr')
        leg_el.remove(pos); leg_el.insert(0, pos)
        # manualLayout positions in chart-fraction units
        layout = ET.Element(qn('c:layout'))
        ml = ET.SubElement(layout, qn('c:manualLayout'))
        ET.SubElement(ml, qn('c:xMode')).set('val', 'edge')
        ET.SubElement(ml, qn('c:yMode')).set('val', 'edge')
        ET.SubElement(ml, qn('c:x')).set('val', legend_pos[0])
        ET.SubElement(ml, qn('c:y')).set('val', legend_pos[1])
        ET.SubElement(ml, qn('c:w')).set('val', legend_pos[2])
        ET.SubElement(ml, qn('c:h')).set('val', legend_pos[3])
        pos.addnext(layout)
        # White fill behind legend
        for old in leg_el.findall(qn('c:spPr')):
            leg_el.remove(old)
        leg_spPr = ET.Element(qn('c:spPr'))
        sf = ET.SubElement(leg_spPr, qn('a:solidFill'))
        clr = ET.SubElement(sf, qn('a:srgbClr')); clr.set('val', 'FFFFFF')
        ln = ET.SubElement(leg_spPr, qn('a:ln')); ln.set('w', '6350')
        lf = ET.SubElement(ln, qn('a:solidFill'))
        lc = ET.SubElement(lf, qn('a:srgbClr')); lc.set('val', '0B2B4E')
        layout.addnext(leg_spPr)
    return chart_shape


def _make_xy_line_chart(slide, x, y, w, h, *, series, x_title, y_title,
                          x_min=0, x_max=None, x_unit=None,
                          y_min=0, y_max=None, y_unit=None,
                          legend=False, legend_pos=None, smooth=False):
    """XY-scatter line+markers chart with the deck's standard styling.

    Use when you need data points to sit at arbitrary X-positions (e.g.,
    plotting MPL at the midpoint of each L-interval) while keeping tick
    marks at standard L-values.  Both axes are numeric value axes.

    series: list of ``(name, [(x, y), ...], color: RGBColor, marker: str)``.
    smooth: True for XY_SCATTER_SMOOTH (cubic spline through points),
            False for XY_SCATTER_LINES (straight segments).
    """
    _add_graphicframe_shadow(slide, x, y, w, h)
    cd = XyChartData()
    for name, points, _color, _marker in series:
        s = cd.add_series(name)
        for px, py in points:
            s.add_data_point(px, py)
    chart_type = (XL_CHART_TYPE.XY_SCATTER_SMOOTH if smooth
                  else XL_CHART_TYPE.XY_SCATTER_LINES)
    chart_shape = slide.shapes.add_chart(
        chart_type,
        int(x), int(y), int(w), int(h), cd,
    )
    chart = chart_shape.chart
    chart.has_title = False
    chart.has_legend = legend

    # Per-series styling
    for idx, ser in enumerate(chart.series):
        name, points, color, marker = series[idx]
        line = ser.format.line
        line.color.rgb = color
        line.width = Pt(2.5)
        ser_xml = ser._element
        clr_hex = f'{color[0]:02X}{color[1]:02X}{color[2]:02X}'
        for old in ser_xml.findall(qn('c:marker')):
            ser_xml.remove(old)
        m = ET.SubElement(ser_xml, qn('c:marker'))
        sym = ET.SubElement(m, qn('c:symbol')); sym.set('val', marker)
        sz_el = ET.SubElement(m, qn('c:size')); sz_el.set('val', '7')
        sp = ET.SubElement(m, qn('c:spPr'))
        fl = ET.SubElement(sp, qn('a:solidFill'))
        rg = ET.SubElement(fl, qn('a:srgbClr')); rg.set('val', clr_hex)
        ln = ET.SubElement(sp, qn('a:ln'))
        lf = ET.SubElement(ln, qn('a:solidFill'))
        lr = ET.SubElement(lf, qn('a:srgbClr')); lr.set('val', clr_hex)
        for sm in ser_xml.findall(qn('c:smooth')):
            ser_xml.remove(sm)
        smooth = ET.SubElement(ser_xml, qn('c:smooth'))
        smooth.set('val', '0')

    # Both axes are value axes in XY scatter; python-pptx returns the
    # X axis through .category_axis (wrapped as a ValueAxis since
    # catAx_lst is empty) and the Y axis through .value_axis.
    x_ax = chart.category_axis
    x_ax.tick_labels.font.name = "Calibri"
    x_ax.tick_labels.font.size = Pt(10)
    x_ax.tick_labels.font.color.rgb = NAVY
    if x_max is not None:
        x_ax.minimum_scale = x_min
        x_ax.maximum_scale = x_max
    if x_unit is not None:
        x_ax.major_unit = x_unit
    x_ax.has_title = True
    x_ax.axis_title.text_frame.text = x_title
    ar = x_ax.axis_title.text_frame.paragraphs[0].runs[0]
    ar.font.name = "Calibri"; ar.font.size = Pt(12)
    ar.font.bold = True; ar.font.italic = True; ar.font.color.rgb = NAVY

    y_ax = chart.value_axis
    y_ax.tick_labels.font.name = "Calibri"
    y_ax.tick_labels.font.size = Pt(10)
    y_ax.tick_labels.font.color.rgb = NAVY
    if y_max is not None:
        y_ax.minimum_scale = y_min
        y_ax.maximum_scale = y_max
    if y_unit is not None:
        y_ax.major_unit = y_unit
    y_ax.has_title = True
    y_ax.axis_title.text_frame.text = y_title
    ar = y_ax.axis_title.text_frame.paragraphs[0].runs[0]
    ar.font.name = "Calibri"; ar.font.size = Pt(12)
    ar.font.bold = True; ar.font.italic = True; ar.font.color.rgb = NAVY

    _add_dashed_gridlines(x_ax._element)
    _add_dashed_gridlines(y_ax._element)

    # Legend
    if legend and legend_pos is not None:
        leg_el = chart.legend._element
        chart.legend.font.name = "Calibri"
        chart.legend.font.size = Pt(11)
        chart.legend.font.color.rgb = NAVY
        chart.legend.include_in_layout = False
        for old in leg_el.findall(qn('c:layout')):
            leg_el.remove(old)
        for old in leg_el.findall(qn('c:legendPos')):
            leg_el.remove(old)
        pos = ET.SubElement(leg_el, qn('c:legendPos')); pos.set('val', 'tr')
        leg_el.remove(pos); leg_el.insert(0, pos)
        layout = ET.Element(qn('c:layout'))
        ml = ET.SubElement(layout, qn('c:manualLayout'))
        ET.SubElement(ml, qn('c:xMode')).set('val', 'edge')
        ET.SubElement(ml, qn('c:yMode')).set('val', 'edge')
        ET.SubElement(ml, qn('c:x')).set('val', legend_pos[0])
        ET.SubElement(ml, qn('c:y')).set('val', legend_pos[1])
        ET.SubElement(ml, qn('c:w')).set('val', legend_pos[2])
        ET.SubElement(ml, qn('c:h')).set('val', legend_pos[3])
        pos.addnext(layout)
        for old in leg_el.findall(qn('c:spPr')):
            leg_el.remove(old)
        leg_spPr = ET.Element(qn('c:spPr'))
        sf = ET.SubElement(leg_spPr, qn('a:solidFill'))
        clr = ET.SubElement(sf, qn('a:srgbClr')); clr.set('val', 'FFFFFF')
        ln = ET.SubElement(leg_spPr, qn('a:ln')); ln.set('w', '6350')
        lf = ET.SubElement(ln, qn('a:solidFill'))
        lc = ET.SubElement(lf, qn('a:srgbClr')); lc.set('val', '0B2B4E')
        layout.addnext(leg_spPr)

    return chart_shape


def _omml_acc_overline(symbol):
    """Inline OMML accent (overline / bar) on a math symbol.

    symbol: e.g. 'K' or 'L' – the variable to wear the bar.
    Returns an OMML fragment intended to be embedded inside an <m:oMath>.
    """
    return (
        '<m:acc>'
          '<m:accPr><m:chr m:val="̅"/></m:accPr>'
          '<m:e>'
            '<m:r>'
              '<a:rPr lang="en-US" b="0" i="1">'
                '<a:latin typeface="Cambria Math"/>'
              '</a:rPr>'
              f'<m:t>{symbol}</m:t>'
            '</m:r>'
          '</m:e>'
        '</m:acc>'
    )


def _add_mixed_textbox(slide, left, top, width, height, segments, *,
                        align=PP_ALIGN.LEFT, default_color=NAVY,
                        default_size=24,
                        margin_left=None, margin_right=None,
                        margin_top=None, margin_bottom=None):
    """Build a textbox whose paragraphs mix plain text runs and inline OMML.

    segments: list of (kind, content, opts) tuples, with kind ∈ {"text",
    "omml", "break"}.  "break" inserts a new paragraph.  Opts may set
    `size`, `bold`, `italic`, `color`, `font` per run.
    """
    box = slide.shapes.add_textbox(int(left), int(top),
                                     int(width), int(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05) if margin_left is None else margin_left
    tf.margin_right = Inches(0.05) if margin_right is None else margin_right
    tf.margin_top = Inches(0.0) if margin_top is None else margin_top
    tf.margin_bottom = Inches(0.0) if margin_bottom is None else margin_bottom

    align_attr = ''
    if align == PP_ALIGN.CENTER: align_attr = ' algn="ctr"'
    elif align == PP_ALIGN.RIGHT: align_attr = ' algn="r"'

    def _start_para():
        return [f'<a:p xmlns:a="{A_NS}" xmlns:m="{M_NS}" xmlns:a14="{A14_NS}">',
                f'<a:pPr{align_attr}/>' if align_attr else '']

    paragraphs = [_start_para()]
    for kind, content, opts in segments:
        if kind == 'break':
            paragraphs[-1].append('<a:endParaRPr lang="en-US"/></a:p>')
            paragraphs.append(_start_para())
            continue
        size_pt = int(opts.get('size', default_size) * 100)
        color = opts.get('color', default_color)
        clr_hex = f'{color[0]:02X}{color[1]:02X}{color[2]:02X}'
        bold_attr = ' b="1"' if opts.get('bold') else ''
        italic_attr = ' i="1"' if opts.get('italic') else ''
        font = opts.get('font', 'Calibri')
        if kind == 'text':
            paragraphs[-1].append(
                f'<a:r><a:rPr lang="en-US" sz="{size_pt}"{bold_attr}{italic_attr}>'
                f'<a:solidFill><a:srgbClr val="{clr_hex}"/></a:solidFill>'
                f'<a:latin typeface="{font}"/>'
                f'</a:rPr><a:t>{content}</a:t></a:r>'
            )
        elif kind == 'omml':
            paragraphs[-1].append(
                f'<a14:m><m:oMath>{content}</m:oMath></a14:m>'
            )
    paragraphs[-1].append('<a:endParaRPr lang="en-US"/></a:p>')
    full_xml = ''.join(''.join(p) for p in paragraphs)

    txBody = tf._txBody
    for old in list(txBody.findall(qn('a:p'))):
        txBody.remove(old)
    # We have to parse each <a:p> separately so the namespaces resolve
    for p_str in full_xml.split('</a:p>'):
        if not p_str.strip(): continue
        p_xml = p_str + '</a:p>'
        new_p = ET.fromstring(p_xml)
        txBody.append(new_p)
        # Apply size+color to any OMML m:r elements inside this paragraph.
        # Color is set to ``default_color`` ONLY when no per-run solidFill
        # is already present — this lets callers tint individual OMML runs
        # via the optional ``color=`` argument on ``_omml_run`` / ``_omml_text``
        # (e.g., the green ΔL / ΔQ in the slide-14 Convention box) without
        # being silently overridden here.
        clr_hex = f'{default_color[0]:02X}{default_color[1]:02X}{default_color[2]:02X}'
        for r in new_p.iter(qn('m:r')):
            arPr = r.find(qn('a:rPr'))
            if arPr is None:
                arPr = ET.Element(qn('a:rPr'))
                arPr.set('lang', 'en-US')
                r.insert(0, arPr)
            arPr.set('sz', str(int(default_size * 100)))
            if arPr.find(qn('a:solidFill')) is None:
                sf = ET.SubElement(arPr, qn('a:solidFill'))
                srgb = ET.SubElement(sf, qn('a:srgbClr'))
                srgb.set('val', clr_hex)
    return box


def _add_math_equation(slide, left, top, width, height, omml_content, *,
                       size_pt=32, color=NAVY, fill=None, line=None):
    """Place an OMML equation in a textbox on the slide.

    omml_content: a string built from _omml_* helpers (without the outer
    <m:oMathPara> wrapper).
    """
    left, top, width, height = int(left), int(top), int(width), int(height)
    box = slide.shapes.add_textbox(left, top, width, height)

    if fill is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
    if line is not None:
        box.line.color.rgb = line
        box.line.width = Pt(0.75)

    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Replace the default <a:p> with one that hosts a math zone (<a14:m>)
    # containing the oMathPara.  The <a14:m> wrapper is REQUIRED by
    # PowerPoint to recognise OMML inside a textbox – without it PPT just
    # shows empty boxes.
    txBody = tf._txBody
    for p in list(txBody.findall(qn('a:p'))):
        txBody.remove(p)

    sz = int(size_pt * 100)
    clr_hex = '{:02X}{:02X}{:02X}'.format(color[0], color[1], color[2])

    # Use unique namespace prefixes; lxml will resolve them when parsing.
    p_xml = (
        f'<a:p xmlns:a="{A_NS}" xmlns:m="{M_NS}" xmlns:a14="{A14_NS}">'
        f'<a:pPr algn="ctr">'
        f'<a:defRPr sz="{sz}" b="0" i="1">'
        f'<a:solidFill><a:srgbClr val="{clr_hex}"/></a:solidFill>'
        f'<a:latin typeface="Cambria Math"/>'
        f'</a:defRPr>'
        f'</a:pPr>'
        f'<a14:m>'
        f'<m:oMathPara>'
        f'<m:oMathParaPr><m:jc m:val="centerGroup"/></m:oMathParaPr>'
        f'<m:oMath>{omml_content}</m:oMath>'
        f'</m:oMathPara>'
        f'</a14:m>'
        f'<a:endParaRPr lang="en-US" sz="{sz}"/>'
        f'</a:p>'
    )

    new_p = ET.fromstring(p_xml)
    txBody.append(new_p)

    # Set size and color on every OMML run's <a:rPr> so the text renders at
    # the right size.  (The defRPr above is the fallback.)
    for r in new_p.iter(qn('m:r')):
        arPr = r.find(qn('a:rPr'))
        if arPr is None:
            arPr = ET.Element(qn('a:rPr'))
            r.insert(0, arPr)
        arPr.set('sz', str(sz))
        if arPr.get('lang') is None:
            arPr.set('lang', 'en-US')
        for sf in arPr.findall(qn('a:solidFill')):
            arPr.remove(sf)
        sf = ET.SubElement(arPr, qn('a:solidFill'))
        srgb = ET.SubElement(sf, qn('a:srgbClr'))
        srgb.set('val', clr_hex)
    return box


# --------------------------------------------------------------------------
# Convenience: build typical formula structures
# --------------------------------------------------------------------------

def _formula_bang_for_buck(op=' = ', reverse=False):
    """Bang-for-the-buck rule with stacked fractions and subscripts.

    Default order:  MP_K / p_K   [op]   MP_L / w
    reverse=True:   MP_L / w     [op]   MP_K / p_K
        (Use when the L-fraction is the larger one, so reading
         left-to-right matches "workers' MP/$ > robots' MP/$".)
    """
    mp_k = _omml_sub(_omml_run('MP'), _omml_run('K'))
    p_k  = _omml_sub(_omml_run('p'),  _omml_run('K'))
    mp_l = _omml_sub(_omml_run('MP'), _omml_run('L'))
    w    = _omml_run('w')
    frac_k = _omml_frac(mp_k, p_k)
    frac_l = _omml_frac(mp_l, w)
    if reverse:
        return frac_l + _omml_text(op) + frac_k
    return frac_k + _omml_text(op) + frac_l


def _formula_optimal_inputs():
    """The combined optimum condition with both fractions."""
    return _formula_bang_for_buck()


def _formula_mp_ratio(input_name='K', price_symbol='p'):
    """MP_X / p_X  (one side of the rule)."""
    base = _omml_sub(_omml_run('MP'), _omml_run(input_name))
    if price_symbol == 'w':
        den = _omml_run('w')
    else:
        den = _omml_sub(_omml_run(price_symbol), _omml_run(input_name))
    return _omml_frac(base, den)


def _add_half_textbox(slide, left, top, width, height, items, *,
                      size=22, line_spacing_pts=14, color=NAVY):
    """A simple half-page text block with bulleted or unbulleted lines."""
    box = slide.shapes.add_textbox(int(left), int(top), int(width), int(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            text, is_bullet = item
        else:
            text, is_bullet = item, True
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        if i > 0:
            pPr = p._p.get_or_add_pPr()
            spcBef = ET.SubElement(pPr, qn('a:spcBef'))
            pts = ET.SubElement(spcBef, qn('a:spcPts'))
            pts.set('val', str(line_spacing_pts * 100))
        run = p.add_run()
        run.text = text
        run.font.name = "Calibri"
        run.font.size = Pt(size)
        run.font.bold = False
        run.font.color.rgb = color
        if is_bullet:
            _set_bullet_char(p, char='▪', color=NAVY,
                              mar_l=342900, indent=-342900, size_pct=100)
    return box


# --------------------------------------------------------------------------
# Slide builders
# --------------------------------------------------------------------------

SECTION_TAG_FRONT = "Module 3 · Front Matter"


def slide_1(prs):
    s = make_title_slide(prs)
    _set_notes(s, (
        "Welcome – this is Module 3, Production and Costs. Last time we "
        "wrapped up the demand side; tonight we tackle the supply side and "
        "how output depends on inputs. By the end you'll have all the "
        "pieces you need for Module 4, where we put demand and costs "
        "together to find profit-maximizing decisions."
    ))


def slide_2(prs):
    bullets = [
        "If your internet permits, keep video on, mic muted",
        ("Class discussions: raise hand", 0),
        ("Will call on you to unmute", 1),
        ("Chat on Zoom: only clarifying questions to TA", 0),
        ("TA will filter questions and raise hand if needed", 1),
        "Group discussions / exercises: breakout rooms",
    ]

    def add_zoom_logo(slide):
        logo_path = OUT_DIR / "_zoom_logo.png"
        if logo_path.exists():
            slide.shapes.add_picture(
                str(logo_path),
                Inches(11.3), Inches(0.65),
                width=Inches(1.6),
            )

    s = make_content_bulleted(
        prs,
        page_num=2,
        section_tag="Module 3 · Logistics",
        title="Zoom-Specific Logistics",
        bullets=bullets,
        size=32, sub_size=28,
        line_spacing_pts=22,
        extras=add_zoom_logo,
    )
    _set_notes(s, (
        "Quick housekeeping before we dive in: Zoom rules. Video on if your "
        "bandwidth allows, mic muted by default. Raise hand to be called on; "
        "use the TA chat for clarifying questions. Group exercises run in "
        "breakout rooms – I'll move between them. Then we move on."
    ))


def slide_announcements(prs):
    """Slide 3 – Announcements (midterm logistics).

    Reintroduced 2026-05-15 to mirror the original deck's slide 3.
    Dates left as ``{{MIDTERM_WINDOW}}`` / ``{{TA_WINDOW}}`` placeholders
    so the actual 2026 dates can be filled in later in PowerPoint.
    """
    bullets = [
        ("Midterm logistics", 0),
        ("3.5-hour window at home, any time during {{MIDTERM_WINDOW}}", 1),
        ("Guaranteed TA availability: {{TA_WINDOW}}", 1),
        ("Material covered", 0),
        ("All material from Modules 1 and 2 (includes PS 1 + 2)", 1),
        ("Problem-solving exercises similar to PS 1 + 2", 1),
        ("Review sessions during the midterm week", 0),
    ]
    s = make_content_bulleted(
        prs,
        page_num=3,
        section_tag="Module 3 · Announcements",
        title="Announcements",
        bullets=bullets,
        size=30, sub_size=26,
        line_spacing_pts=12,
    )
    _set_notes(s, (
        "Before we dive into Module 3 – two quick announcements about the "
        "midterm. It's a 6-hour at-home window over the dates shown; pick "
        "any contiguous block within that window. We'll have a TA on call "
        "for one guaranteed support window – use it if you'd like to ask "
        "live questions. Material covers Modules 1 and 2, including both "
        "problem sets. Review sessions will run during the week of the "
        "midterm; details by email."
    ))


def slide_3(prs):
    bullets = [
        ("The law of demand", 0),
        ("Holding everything else constant, if P falls, Q rises", 1),
        ("Elasticities", 0),
        ("Responsiveness of demand to own-price, income, competitors' price", 1),
        ("Demand and revenue", 0),
        ("Own-price elasticity drives total revenue; marginal revenue via the 3-step method", 1),
        ("Demand estimation", 0),
        ("Market experimentation and regression analysis", 1),
    ]
    s = make_content_bulleted(
        prs,
        page_num=4,
        section_tag="Module 3 · Recap",
        title="Recap of Module 2",
        bullets=bullets,
        size=32, sub_size=26,
        line_spacing_pts=10,
    )
    _set_notes(s, (
        "A 60-second reminder of where Module 2 left us: demand curves, "
        "price elasticity, and marginal revenue. The revenue side is settled; "
        "tonight we crack open the cost side. Once both sides are in hand, "
        "profit-maximization in Module 4 falls out almost mechanically."
    ))


def slide_4(prs):
    """Course-roadmap flowchart preserving the original 4-module structure.

    Layout (mirroring the source slide, beautified in template colors):

        ┌───────────────────────────────────────────┐
        │ 1. Basic Principles and Economic Way ...  │  (faded)
        └─────────────────┬─────────────────────────┘
              ┌───────────┴──────────────┐
              ▼                          ▼
        ┌──────────────────┐    ┌──────────────────┐
        │ 2. Value & Demand│    │ 3. Supply & Cost │  (←navy, current)
        └─────────┬────────┘    └────────┬─────────┘   ←┐ "You are here"
                  └─────────┬────────────┘             │
                            ▼                          │
                  ┌───────────────────────────┐        │
                  │ 4. Markets, Pricing, ...  │  (faded)
                  └───────────────────────────┘

    The current module (3) is highlighted in navy with a gold "you are here"
    arrow; past/future modules render in faded grey.
    """

    def draw(slide):
        # Geometry — taller boxes to accommodate larger type
        box_h = Inches(0.85)
        narrow_w = Inches(4.6)
        wide_w = Inches(8.6)
        gap = Inches(0.3)

        slide_mid = SLIDE_W // 2

        # Row 1 – top module (faded)
        top_x = slide_mid - wide_w // 2
        top_y = Inches(2.0)
        _add_rounded_filled_box(slide, top_x, top_y, wide_w, box_h,
                                  "1. Basic Principles and Economic Way of Thinking",
                                  fill=FADED, text_color=WHITE, size=24, bold=True)

        # Row 2 – two parallel modules
        row2_y = Inches(3.65)
        left_x = slide_mid - gap // 2 - narrow_w
        right_x = slide_mid + gap // 2
        _add_rounded_filled_box(slide, left_x, row2_y, narrow_w, box_h,
                                  "2. Value and Demand",
                                  fill=FADED, text_color=WHITE, size=26, bold=True)
        # Current module (navy)
        _add_rounded_filled_box(slide, right_x, row2_y, narrow_w, box_h,
                                  "3. Supply and Cost",
                                  fill=NAVY, text_color=WHITE, size=26, bold=True)

        # Row 3 – bottom module (faded)
        bot_x = slide_mid - wide_w // 2
        bot_y = Inches(5.5)
        _add_rounded_filled_box(slide, bot_x, bot_y, wide_w, box_h,
                                  "4. Markets, Pricing, and Strategy",
                                  fill=FADED, text_color=WHITE, size=24, bold=True)

        # Connectors — top down to row 2 (faded grey lines).  Thicker
        # than before per user request 2026-05-15.
        top_bottom_y = top_y + box_h
        row2_top_y = row2_y
        _add_arrow(slide,
                    (top_x + wide_w // 2, top_bottom_y),
                    (left_x + narrow_w // 2, row2_top_y),
                    color=FADED, weight_pt=3.0, head=True)
        _add_arrow(slide,
                    (top_x + wide_w // 2, top_bottom_y),
                    (right_x + narrow_w // 2, row2_top_y),
                    color=NAVY, weight_pt=3.5, head=True)

        # Row 2 down to row 3
        row2_bottom_y = row2_y + box_h
        row3_top_y = bot_y
        _add_arrow(slide,
                    (left_x + narrow_w // 2, row2_bottom_y),
                    (bot_x + wide_w // 2, row3_top_y),
                    color=FADED, weight_pt=3.0, head=True)
        _add_arrow(slide,
                    (right_x + narrow_w // 2, row2_bottom_y),
                    (bot_x + wide_w // 2, row3_top_y),
                    color=FADED, weight_pt=3.0, head=True)

        # "We are here" — gold UP-pointing arrow positioned BELOW box 3,
        # head pointing INTO the bottom edge of the box.  Vertical budget
        # is tight: ~1.0" of clearance between row 2 bottom (4.50") and
        # row 3 top (5.50"), so the arrow + label must fit inside that.
        # Shifted right by ~0.55" per user request 2026-05-15.
        right_shift = Inches(0.55)
        arrow_w = Inches(0.55)
        arrow_h = Inches(0.55)
        arrow_left = right_x + (narrow_w - arrow_w) // 2 + right_shift
        arrow_top = row2_y + box_h + Inches(0.05)
        _add_arrow_shape(slide, arrow_left, arrow_top, arrow_w, arrow_h,
                          direction="up", fill=GOLD)
        # Label directly below the arrow, shifted with the arrow
        _add_text(slide, right_x + right_shift,
                   arrow_top + arrow_h + Inches(0.02),
                   narrow_w, Inches(0.32),
                   "we are here", size=16, italic=True, bold=True,
                   color=GOLD, font="Calibri", align=PP_ALIGN.CENTER)

    s = make_diagram_slide(
        prs, page_num=5,
        section_tag="Module 3 · Course Roadmap",
        title="Agenda for the Class",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "This is where Module 3 fits in the 405 course as a whole. We "
        "started with the Basic Principles and the Economic Way of Thinking. "
        "Module 2 covered Value and Demand – the customer side. Tonight is "
        "Module 3, Supply and Cost – the firm side. Module 4 ties it all "
        "together in Markets, Pricing, and Strategy. So tonight you're "
        "filling in the right-hand box on this map."
    ))


def slide_5(prs):
    """'Big Picture' diagram – Production Functions + Costs (Module 3) feed
    into Output Decisions (Module 4); Demand (Module 2) also feeds in.

    Layout:

        ┌─────────────────────┐         ┌──────────────┐
        │ Production Functions│         │   Demand     │  (faded, M2)
        │  inputs → output    │         │   (Module 2) │
        └─────────────────────┘         └──────┬───────┘
                  │                            │
                  ▼                            │
        ┌─────────────────────┐                │
        │       Costs         │                │
        │ what each unit costs│                │
        └─────────────────────┘                │
                  │                            │
                  └────────────┬───────────────┘
                               ▼
                  ┌───────────────────────────┐
                  │   Output Decisions        │  (gold, M4 preview)
                  │   (Module 4: pricing &    │
                  │    profit-maximization)   │
                  └───────────────────────────┘
    """

    def draw(slide):
        # Geometry: 3-column grid, left (M3 stack), right (M2), bottom-center (M4)
        # Box wording restored to the original deck's slide 8 (single-line:
        # no parenthetical sub-text), so heights are shorter than before.
        box_h = Inches(1.05)
        small_h = Inches(1.05)
        m3_w = Inches(4.8)
        m2_w = Inches(4.2)
        m4_w = Inches(7.0)

        # Module 3 left stack
        m3_x = Inches(0.5)
        prod_y = Inches(2.05)
        costs_y = Inches(3.85)
        _add_rounded_filled_box(slide, m3_x, prod_y, m3_w, box_h,
                                  "Production Functions",
                                  fill=NAVY, text_color=WHITE,
                                  size=36, bold=True)
        _add_rounded_filled_box(slide, m3_x, costs_y, m3_w, box_h,
                                  "Costs",
                                  fill=NAVY, text_color=WHITE,
                                  size=36, bold=True)

        # Module 2 right (faded — already covered)
        m2_x = Inches(8.6)
        m2_y = Inches(2.05)
        _add_rounded_filled_box(slide, m2_x, m2_y, m2_w, small_h,
                                  "Demand",
                                  fill=FADED, text_color=WHITE,
                                  size=36, bold=True)

        # Module 4 bottom-center (gold — coming up)
        m4_x = (SLIDE_W - m4_w) // 2
        m4_y = Inches(5.8)
        _add_rounded_filled_box(slide, m4_x, m4_y, m4_w, box_h,
                                  "Output Decisions",
                                  fill=GOLD, text_color=WHITE,
                                  size=36, bold=True)

        # Arrows — thicker than before per user request 2026-05-15.
        # Production → Costs (vertical, inside the M3 stack)
        _add_arrow(slide,
                    (m3_x + m3_w // 2, prod_y + box_h),
                    (m3_x + m3_w // 2, costs_y),
                    color=NAVY, weight_pt=3.5, head=True)
        # Costs → Output Decisions (diagonal down to centre)
        _add_arrow(slide,
                    (m3_x + m3_w // 2, costs_y + box_h),
                    (int(m4_x + m4_w * 0.30), m4_y),
                    color=NAVY, weight_pt=3.5, head=True)
        # Demand → Output Decisions (diagonal down to centre-right)
        _add_arrow(slide,
                    (m2_x + m2_w // 2, m2_y + small_h),
                    (int(m4_x + m4_w * 0.70), m4_y),
                    color=FADED, weight_pt=3.0, head=True)

    s = make_diagram_slide(
        prs, page_num=6,
        section_tag="Module 3 · Production · Big Picture",
        title="Big Picture of Module 3",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "The big-picture frame for the night. Every executive decision "
        "you've ever made – pricing, hiring, capacity, sourcing, "
        "outsourcing – is at some level a production-and-cost decision. "
        "Tonight builds the left-hand chain in this diagram: Production "
        "Functions describe how inputs become output, and Costs translate "
        "those inputs into dollars. Combined with Demand from Module 2, "
        "they give us Output Decisions in Module 4 – pricing and profit-"
        "maximization. Click through the boxes one at a time to walk the "
        "students through the flow."
    ))


# --------------------------------------------------------------------------
# Batch 2 – Part 1 (Production) §1.1 Short Run
# (Old textual slide_6 outline dropped – replaced by slide_concept_map.)
# --------------------------------------------------------------------------

SECTION_TAG_P1 = "Module 3 · Production · Short Run"
SECTION_TAG_P1_DIV = "Module 3 · Agenda"


# --------------------------------------------------------------------------
# Slide 7 (NEW): Concept-map / network graph
# Visualises how the components from the Outline (slide 6) relate.
# --------------------------------------------------------------------------

def slide_concept_map(prs):
    """Two parallel trees showing how Module 3's pieces fit together.

    Production rules use OMML (Cambria Math) for proper TeX-style
    formula rendering with real subscripts and stacked fractions.

    Visual story for EMBA pacing:
      • Production function → two time horizons → two decision rules
      • Each rule IS the optimization condition in its time horizon
      • Following these rules minimizes cost for a given Q  →  THAT is
        the bridge to the Costs side, which describes those costs
      • A small gold callout flags the cleanest MB = MC example
        (short-run hiring); bang-for-the-buck is a related MB / $ rule
    """

    def _formula_box(slide, x, y, w, h, label, omml_expr,
                      *, label_size=12, formula_size=22):
        """Navy filled box with a Calibri label on top + OMML formula below."""
        _add_filled_box(slide, x, y, w, h, "", fill=NAVY)
        _add_text(slide, x, y + Inches(0.06), w, Inches(0.30),
                   label, size=label_size, bold=True, color=WHITE,
                   font="Calibri", align=PP_ALIGN.CENTER)
        _add_math_equation(
            slide, x, y + Inches(0.36), w, h - Inches(0.42),
            omml_expr, size_pt=formula_size, color=WHITE,
        )

    def draw(slide):
        # Cluster header labels (section headers for the two halves)
        _add_text(slide, Inches(0.3), Inches(1.35), Inches(6.0), Inches(0.45),
                   "PRODUCTION", size=22, bold=True, color=GRAY,
                   font="Calibri", align=PP_ALIGN.CENTER)
        _add_text(slide, Inches(7.0), Inches(1.35), Inches(6.0), Inches(0.45),
                   "COSTS", size=22, bold=True, color=GRAY,
                   font="Calibri", align=PP_ALIGN.CENTER)

        # ---- PRODUCTION cluster (left) ---------------------------------

        # N1: Production function — label on top + OMML Q = f(K, L)
        n1_x, n1_y = Inches(1.0), Inches(1.90)
        n1_w, n1_h = Inches(4.6), Inches(0.95)
        eq_pf = (
            _omml_run('Q') + _omml_text(' = ') + _omml_run('f') +
            _omml_text('(') + _omml_run('K') + _omml_text(', ') +
            _omml_run('L') + _omml_text(')')
        )
        _formula_box(slide, n1_x, n1_y, n1_w, n1_h,
                      "Production function", eq_pf,
                      label_size=16, formula_size=24)

        # N2 / N3: Short Run | Long Run
        sr_x, lr_x = Inches(0.3), Inches(3.5)
        sr_y = Inches(3.05)
        col_w, col_h = Inches(2.8), Inches(0.65)
        _add_filled_box(slide, sr_x, sr_y, col_w, col_h,
                         "Short Run\n(K fixed, L flexible)",
                         fill=NAVY, text_color=WHITE, size=14, bold=True)
        _add_filled_box(slide, lr_x, sr_y, col_w, col_h,
                         "Long Run\n(Both K and L flexible)",
                         fill=NAVY, text_color=WHITE, size=14, bold=True)

        # N4 / N5: decision rules with OMML formulas
        rule_y, rule_h = Inches(3.90), Inches(1.05)

        # Hire until  MRPL = w
        eq_hire = (
            _omml_text('MRPL') + _omml_text(' = ') + _omml_run('w')
        )
        _formula_box(slide, sr_x, rule_y, col_w, rule_h,
                      "Hire until", eq_hire,
                      label_size=12, formula_size=24)

        # Bang-for-the-buck:  MP_K / p_K  =  MP_L / w
        mp_k = _omml_sub(_omml_text('MP'), _omml_run('K'))
        p_k  = _omml_sub(_omml_run('p'),  _omml_run('K'))
        mp_l = _omml_sub(_omml_text('MP'), _omml_run('L'))
        wvar = _omml_run('w')
        eq_bfb = (
            _omml_frac(mp_k, p_k) + _omml_text(' = ') + _omml_frac(mp_l, wvar)
        )
        _formula_box(slide, lr_x, rule_y, col_w, rule_h,
                      "Bang-for-the-buck", eq_bfb,
                      label_size=12, formula_size=17)

        # Small italic labels under each rule — tying the rule to
        # optimization in its time horizon
        opt_y = rule_y + rule_h + Inches(0.05)
        _add_text(slide, sr_x, opt_y, col_w, Inches(0.28),
                   "↑  short-run optimization",
                   size=12, italic=True, color=GOLD, bold=True,
                   font="Calibri", align=PP_ALIGN.CENTER)
        _add_text(slide, lr_x, opt_y, col_w, Inches(0.28),
                   "↑  long-run optimization",
                   size=12, italic=True, color=GOLD, bold=True,
                   font="Calibri", align=PP_ALIGN.CENTER)

        # ---- COSTS cluster (right) -------------------------------------

        c1_x, c1_y = Inches(7.4), Inches(1.90)
        c1_w, c1_h = Inches(5.6), Inches(0.95)
        # Cost-types header: label on top + OMML acronyms below in
        # upright Cambria Math (matches the deck's TeX-style convention).
        # Six acronyms drawn from the cost-concepts section: total fixed,
        # total variable, average fixed, average variable, average total,
        # and marginal cost.
        eq_costs = (
            _omml_text('TFC') + _omml_text('  /  ') +
            _omml_text('TVC') + _omml_text('  /  ') +
            _omml_text('AFC') + _omml_text('  /  ') +
            _omml_text('AVC') + _omml_text('  /  ') +
            _omml_text('ATC') + _omml_text('  /  ') +
            _omml_text('MC')
        )
        _formula_box(slide, c1_x, c1_y, c1_w, c1_h,
                      "Cost types", eq_costs,
                      label_size=16, formula_size=20)

        # Three parallel children of C1 — Fixed / Marginal / Average costs.
        # Each is a key decision-relevant view of costs.
        child_y, child_h = Inches(3.05), Inches(1.90)
        gap = Inches(0.2)
        cw = (c1_w - 2 * gap) // 3
        c2_x = c1_x
        c3_x = c1_x + cw + gap
        c4_x = c1_x + 2 * (cw + gap)
        _add_filled_box(slide, c2_x, child_y, cw, child_h,
                         "Fixed costs\n\n⇒ ignore if\nthey are sunk",
                         fill=NAVY, text_color=WHITE, size=14, bold=True)
        _add_filled_box(slide, c3_x, child_y, cw, child_h,
                         "Marginal Costs\n\ncrucial for optimal\noutput decisions:\nMR = MC\n(In Module 4)",
                         fill=NAVY, text_color=WHITE, size=13, bold=True)
        _add_filled_box(slide, c4_x, child_y, cw, child_h,
                         "Average Costs\n\nmostly for\naccounting purposes",
                         fill=NAVY, text_color=WHITE, size=13, bold=True)

        # ---- Within-cluster arrows (NAVY) ------------------------------

        # Production
        _add_arrow(slide,
                    (n1_x + n1_w // 2, n1_y + n1_h),
                    (sr_x + col_w // 2, sr_y),
                    color=NAVY, weight_pt=2.5, head=True)
        _add_arrow(slide,
                    (n1_x + n1_w // 2, n1_y + n1_h),
                    (lr_x + col_w // 2, sr_y),
                    color=NAVY, weight_pt=2.5, head=True)
        _add_arrow(slide,
                    (sr_x + col_w // 2, sr_y + col_h),
                    (sr_x + col_w // 2, rule_y),
                    color=NAVY, weight_pt=2.5, head=True)
        _add_arrow(slide,
                    (lr_x + col_w // 2, sr_y + col_h),
                    (lr_x + col_w // 2, rule_y),
                    color=NAVY, weight_pt=2.5, head=True)

        # Costs – C1 fans out into the three parallel children
        for cx in (c2_x, c3_x, c4_x):
            _add_arrow(slide,
                        (c1_x + c1_w // 2, c1_y + c1_h),
                        (cx + cw // 2, child_y),
                        color=NAVY, weight_pt=2.5, head=True)

        # ---- MB = MC anchor (12-point star, distinctive shape) ---------

        # The MB = MC anchor uses a 12-point star (continuous outline) so
        # it stands out from every other rectangular box on the slide.
        # The same shape is reused on every slide where MB = MC appears.
        sun_w, sun_h = Inches(1.8), Inches(1.35)
        sun_x = sr_x + (col_w - sun_w) // 2          # centred under MRPL = w
        sun_y = Inches(5.45)
        _add_anchor_burst(
            slide, sun_x, sun_y, sun_w, sun_h,
            top_text="MB = MC",
            bottom_text="(of labor)",
            top_size=16, bottom_size=12,
        )
        # Arrow from burst UP to MRPL = w (the rule where MB = MC lives)
        _add_arrow(slide,
                    (sun_x + sun_w // 2, sun_y),
                    (sr_x + col_w // 2, rule_y + rule_h),
                    color=GOLD, weight_pt=2.5, head=True)

        # ---- Min-cost bridge: ONE inflow (long-run optimization) -------

        # Centred under Bang-for-the-buck (the long-run optimization rule
        # is the cleanest cost-minimization condition – it chooses BOTH
        # K and L for a given Q).
        bridge_x = lr_x - Inches(0.1)
        bridge_y = Inches(5.60)
        bridge_w = Inches(3.0)
        bridge_h = Inches(1.05)
        _add_outlined_box(
            slide, bridge_x, bridge_y, bridge_w, bridge_h,
            "Minimum cost\nfor any given Q",
            fill=WHITE, line=GOLD, text_color=NAVY,
            size=18, bold=True, line_w=2.0,
        )

        # Inflow arrow: Bang-for-the-buck rule  →  Bridge
        # (long-run optimization gives the minimum cost for any Q)
        _add_arrow(slide,
                    (lr_x + col_w // 2, rule_y + rule_h + Inches(0.32)),
                    (bridge_x + bridge_w // 2, bridge_y),
                    color=GOLD, weight_pt=3.0, head=True)

        # ---- Outflow: bridge → cost cluster ----------------------------

        # Arrow departs from the "Minimum cost for any Q" bridge and
        # points up-right at the cost cluster (lands inside the cluster's
        # bottom-left area, NOT into any specific cost child).
        arrow_start = (bridge_x + bridge_w - Inches(0.3),
                        bridge_y + Inches(0.15))
        arrow_end = (c1_x + Inches(0.4), child_y + child_h)
        _add_arrow(slide, arrow_start, arrow_end,
                    color=GOLD, weight_pt=3.0, head=True)

        # Compact label sitting slightly to the right of the outflow
        # arrow (offset so the arrow line is no longer covered by text).
        _add_text(slide, Inches(7.10), Inches(5.45),
                   Inches(1.70), Inches(0.80),
                   "Optimized production leads to minimum-possible costs",
                   size=14, italic=True, bold=True, color=GOLD,
                   font="Calibri", align=PP_ALIGN.CENTER)

        # ---- Second MB = MC star (anchored under Marginal Costs) -------

        # The "MR = MC" rule inside the Marginal Costs box is a specific
        # instance of MB = MC (here marginal benefit is marginal revenue,
        # over OUTPUT choice). Same star pattern as the production-side
        # anchor, keeping the visual convention consistent.
        mc_sun_w, mc_sun_h = Inches(1.8), Inches(1.35)
        mc_sun_x = c3_x + (cw - mc_sun_w) // 2
        mc_sun_y = Inches(5.45)
        _add_anchor_burst(
            slide, mc_sun_x, mc_sun_y, mc_sun_w, mc_sun_h,
            top_text="MB = MC",
            bottom_text="(of output)",
            top_size=16, bottom_size=12,
        )
        # Arrow from this star UP to the Marginal Costs box (Box 2)
        _add_arrow(slide,
                    (mc_sun_x + mc_sun_w // 2, mc_sun_y),
                    (c3_x + cw // 2, child_y + child_h),
                    color=GOLD, weight_pt=2.5, head=True)

        # ---- Scale-implication annotation under Average Costs ----------

        # A small gold-outlined box that names the long-run AC-falls-with-Q
        # condition for economies of scale.
        scale_y = Inches(5.55)
        scale_h = Inches(1.05)
        _add_outlined_box(
            slide, c4_x, scale_y, cw, scale_h,
            "If long-run AC\nfalls with Q\n⇒ Economies\nof Scale",
            fill=WHITE, line=GOLD, text_color=NAVY,
            size=12, bold=True, line_w=1.5,
        )
        # Arrow: Average Costs box DOWN to the Economies-of-Scale annotation
        _add_arrow(slide,
                    (c4_x + cw // 2, child_y + child_h),
                    (c4_x + cw // 2, scale_y),
                    color=NAVY, weight_pt=2.5, head=True)

    s = make_diagram_slide(
        prs, page_num=7,
        section_tag="Module 3 · Concept Map",
        title="How the Pieces of Module 3 Connect",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "Two parallel trees – Production and Costs. On the Production "
        "side, the production function Q = f(K, L, M) branches into two "
        "time horizons (short run, long run), each with its decision "
        "rule rendered in proper math notation: MRPL = w for short-run "
        "hiring and MP_K/p_K = MP_L/w (bang-for-the-buck) for long-run "
        "input mix. The italic gold labels under each rule remind "
        "students that these rules ARE the optimization conditions in "
        "their time horizons. The gold callout flags the cleanest "
        "MB = MC instance (hiring); bang-for-the-buck is the related "
        "marginal-benefit-per-dollar rule, so the callout doesn't "
        "extend there. The gold-outlined bridge band at the bottom "
        "names the link to Costs: solving the optimization conditions "
        "minimizes cost for any given output, and the cost concepts on "
        "the right describe what those minimized costs look like. Use "
        "this slide as the roadmap to return to at section transitions."
    ))


def slide_7(prs):
    """Section divider – Part 1: Production. Layout 2 (Agenda) with Part 1
    highlighted in navy, Part 2 faded grey."""
    s = make_section_agenda(
        prs, page_num=8,
        current_part_idx=0,
        section_tag=SECTION_TAG_P1_DIV,
        title="Part 1: Production – Picking the Right Inputs",
    )
    _set_notes(s, (
        "Entering Part 1 – Production. The core question for the next "
        "40 minutes: how does output depend on inputs? We'll do short-run "
        "hiring decisions first, then long-run input choice."
    ))


def slide_8(prs):
    """The Production Function: Q = f(K, L, etc)."""
    def draw(slide):
        # Big equation, centred near the top of the body region
        _add_text(slide, MARGIN, Inches(2.0), RULE_W, Inches(1.0),
                   "Q = f (K, L, etc)",
                   size=54, bold=True, color=NAVY, font="Calibri",
                   align=PP_ALIGN.CENTER)
        # Variable legend on the LEFT
        legend = [
            ("Q  =  Output", 0),
            ("f   =  a function of inputs:", 0),
            ("K  =  Capital  (physical: factories, machinery, software, IP)", 1),
            ("L  =  Labor", 1),
            ('"etc" can be raw materials, energy...', 1),
        ]
        _add_hierarchical_bullets(
            slide,
            left=MARGIN + Inches(0.5),
            top=Inches(3.3),
            width=Inches(8.5),
            height=Inches(2.6),
            items=legend,
            size=24, sub_size=20, line_spacing_pts=8,
        )
        # The Karl Marx book / "Das Kapital" – source slide had this image at
        # (5.98, 3.13) 1.82x1.82.  Restore it on the right of the legend.
        # User decision: drop the drop-shadow here (book cover is a special
        # case — looks weird with shadow against its existing background).
        _add_source_image(slide, 8, "rId4",
                           left=Inches(10.3), top=Inches(3.3),
                           width=Inches(2.4),
                           shadow=False)
        _add_text(slide, Inches(10.3), Inches(5.85), Inches(2.4), Inches(0.25),
                   "Marx, Das Kapital  (1867)",
                   size=12, italic=True, color=GRAY,
                   align=PP_ALIGN.CENTER, font="Calibri")
        # Bottom explanation callout — Convention-style box (cream-fill
        # rounded rect with navy border).  Pattern documented in the
        # Teaching CLAUDE.md as the preferred format for concept-
        # explanation textboxes; mirrors the "Convention" box on slide 14.
        box_w = Inches(11.0)
        box_h = Inches(1.00)
        box_x = (SLIDE_W - box_w) // 2
        box_y = Inches(6.10)
        _add_convention_box(
            slide, box_x, box_y, box_w, box_h,
            runs=[
                ("A production function transforms inputs into outputs.",
                 {'size': 20, 'bold': True, 'color': NAVY}),
                ("The more efficient this process, the higher is productivity",
                 {'size': 20, 'color': NAVY, 'newline': True}),
            ],
            size=20, align=PP_ALIGN.CENTER,
        )

    s = make_diagram_slide(
        prs, page_num=9,
        section_tag=SECTION_TAG_P1,
        title="The Production Function",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "The production function in its most basic form: output Q is a "
        "function of capital K, labor L, and any other inputs the firm "
        "uses – raw materials, energy, and so on. Capital is more than "
        "buildings – it includes machinery, software, IP, AI systems – "
        "anything you've already paid for that keeps producing. The "
        "callout at the bottom captures the big idea: a production "
        "function maps inputs into outputs, and the more efficiently it "
        "does so, the higher is productivity. Everything else in the "
        "module is built on this expression."
    ))


def slide_9(prs):
    """In the short run, you're stuck with your capacity."""
    bullets = [
        "In the short run, some inputs (fixed factors) cannot be increased or decreased",
        "The long run is a period long enough for all inputs to be variable",
    ]

    def draw_pictures(slide):
        # Picture captions sit ABOVE the pictures (matches the original
        # slide 11's layout); attribution under the LEFT image stays
        # below the picture.
        CAP_TOP = Inches(3.30)
        PIC_TOP = Inches(3.65)
        PIC_W = Inches(5.0)
        PIC_H = Inches(3.0)
        # Captions — wording from the original deck's slide 11
        _add_text(slide, Inches(1.0), int(CAP_TOP), Inches(5.0), Inches(0.32),
                   "Capital fixed in short run",
                   size=14, italic=True, bold=True, color=NAVY,
                   align=PP_ALIGN.CENTER, font="Calibri")
        _add_text(slide, Inches(7.3), int(CAP_TOP), Inches(5.0), Inches(0.32),
                   "Labor (ophthalmologists) fixed in short run",
                   size=14, italic=True, bold=True, color=NAVY,
                   align=PP_ALIGN.CENTER, font="Calibri")
        # LEFT image: Rivian Normal IL assembly-plant floor (CC BY-SA,
        # Wikimedia).  Replaces the stale Tesla factory floor.
        rivian_plant = OUT_DIR / "_rivian_plant.jpg"
        if rivian_plant.exists():
            pic_left = slide.shapes.add_picture(
                str(rivian_plant),
                int(Inches(1.0)), int(PIC_TOP),
                width=int(PIC_W), height=int(PIC_H),
            )
            _apply_picture_style(pic_left)
        # RIGHT image: ophthalmologist photo from the original deck —
        # illustrates skilled labor as a short-run fixed factor (years of
        # specialised training mean the supply can't be ramped on demand).
        pic_right = _add_source_image(slide, 9, "rId5",
                           left=Inches(7.3), top=PIC_TOP,
                           width=PIC_W, height=PIC_H)
        if pic_right is not None:
            _apply_picture_style(pic_right)
        # Tiny attribution under the LEFT image (CC BY-SA author + license).
        attr_top = PIC_TOP + PIC_H + Inches(0.08)
        _add_text(slide, Inches(1.0), int(attr_top),
                   Inches(5.0), Inches(0.18),
                   "Rivian Normal, IL plant  (CC BY-SA, Wikimedia)",
                   size=9, italic=True, color=GRAY,
                   align=PP_ALIGN.CENTER, font="Calibri")

    s = make_content_bulleted(
        prs, page_num=10,
        section_tag=SECTION_TAG_P1,
        title="Short vs. Long Run:  A Critical Distinction",
        bullets=bullets,
        size=28, sub_size=22, line_spacing_pts=14,
        extras=draw_pictures,
    )
    # Shrink the bullet box so it doesn't overlap the images
    _set_notes(s, (
        "The single most important time-scale distinction in this course. "
        "Short run = some inputs are fixed factors — you cannot vary them "
        "within the planning horizon. Long run = a horizon long enough "
        "that EVERY input becomes variable. The two photos give parallel "
        "examples of short-run fixed factors. Left: Rivian's Normal, "
        "Illinois assembly plant — capital (the building, the robot fleet) "
        "is fixed in any given quarter; they can ramp shifts and headcount "
        "but not the four walls. Right: an ophthalmologist — skilled labor "
        "with many years of training is also a fixed factor in the short "
        "run; a hospital cannot conjure up another ophthalmologist on a "
        "month's notice. Both illustrate the same lesson: in the short "
        "run, you optimise around the factors you cannot change."
    ))


# --------------------------------------------------------------------------
# Shared production-function table used by slides 10 (table) and 11 (chart).
# Cobb-Douglas Q = 0.5 · √(K · L)  (i.e., A=0.5, α=β=0.5 – constant returns
# to scale, but BOTH MPK and MPL strictly diminishing).
#
# K and L grids match the ORIGINAL Tesla-Gigafactory table (K = 100 / 200 /
# 300 / 400 robots; L = 0 … 10,000 workers in steps of 1,000).
#
# Integer-rounded values from this exact formula give strictly diminishing
# MPL down every column AND strictly diminishing MPK along every row –
# fixing the bug the original deck had a "CORRECTION" slide about (where
# rounding had inadvertently produced increasing marginal returns).
# --------------------------------------------------------------------------

# --------------------------------------------------------------------------
# Rivian Georgia plant cost data (from Background Material/Module 3 - Make
# vs Buy.xlsx).  Quadratic cost function:  TC = TFC + 200 · Q²,  TFC = $800K.
# Q grid:  10, 20, …, 110 vehicles per week.  Drives the native cost charts
# on slides 55, 56, 57 so all three are consistent and locked to one source.
# --------------------------------------------------------------------------

COST_TFC = 800_000
COST_VAR_COEF = 200
COST_Q_VALS = list(range(10, 111, 10))

def _cost_tc(Q):  return COST_TFC + COST_VAR_COEF * Q * Q
def _cost_tvc(Q): return COST_VAR_COEF * Q * Q
def _cost_avc(Q): return COST_VAR_COEF * Q
def _cost_atc(Q): return _cost_tc(Q) / Q
def _cost_mc(Q, dQ=10):
    return (_cost_tc(Q + dQ) - _cost_tc(Q)) / dQ


# --------------------------------------------------------------------------
# Cross-reference anchors — 0-indexed positions in the built deck.
# When the deck order changes, update these to match new positions; all
# hyperlinks that target these anchors update automatically on rebuild.
# --------------------------------------------------------------------------

SLIDE_IDX_PF_TABLE = 10        # Slide 11: "Rivian's Production Function: R1 Line Weekly Output"
SLIDE_IDX_MPL_CONVENTION = 13  # Slide 14: "Marginal Product of Labor (MPL): Calculation" (Convention callout)


def _add_slide_link_in_slide(slide, search_text, target_slide_idx, *, prs):
    """Post-process: find ``search_text`` in any text run on ``slide``,
    split that run into three runs [before, link, after], and make the
    middle (link) run a hyperlink that jumps to ``prs.slides[target_slide_idx]``
    when clicked.

    Skips runs that already carry an ``<a:hlinkClick>`` — this lets callers
    invoke the helper repeatedly on a slide where the anchor word ("link")
    appears more than once: each call converts the next un-linked
    occurrence, in document order.

    The link run is styled blue + underlined for visual distinction; all
    other run styling (font, size, italic/bold) is preserved from the
    surrounding text.
    """
    target_part = prs.slides[target_slide_idx].part
    rId = slide.part.relate_to(target_part, RT.SLIDE)
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for p in shape.text_frame.paragraphs:
            for run in list(p.runs):
                if search_text not in (run.text or ''):
                    continue
                # Skip runs that are ALREADY a hyperlink (so the next call
                # naturally picks the next un-linked "link" occurrence).
                existing_rPr = run._r.find(qn('a:rPr'))
                if existing_rPr is not None and existing_rPr.find(qn('a:hlinkClick')) is not None:
                    continue

                full = run.text
                idx = full.find(search_text)
                before = full[:idx]
                after = full[idx + len(search_text):]

                orig_r = run._r
                orig_rPr = orig_r.find(qn('a:rPr'))

                # 1. Truncate current run to "before"
                run.text = before

                # 2. Build LINK run (carries the hyperlink)
                link_r = ET.Element(qn('a:r'))
                if orig_rPr is not None:
                    link_rPr = copy.deepcopy(orig_rPr)
                else:
                    link_rPr = ET.SubElement(link_r, qn('a:rPr'))
                # Strip any pre-existing hyperlink, solidFill, and underline
                # from the cloned rPr so our hyperlink styling wins.
                for tag in ('a:hlinkClick', 'a:solidFill'):
                    for old in link_rPr.findall(qn(tag)):
                        link_rPr.remove(old)
                # Underline for visual hyperlink cue
                link_rPr.set('u', 'sng')
                # Blue hyperlink color (PowerPoint's classic hyperlink hue)
                blue = ET.SubElement(link_rPr, qn('a:solidFill'))
                srgb = ET.SubElement(blue, qn('a:srgbClr'))
                srgb.set('val', '0563C1')
                # Hyperlink relationship
                hl = ET.SubElement(link_rPr, qn('a:hlinkClick'))
                hl.set(qn('r:id'), rId)
                hl.set('action', 'ppaction://hlinksldjump')
                if link_rPr not in list(link_r):
                    link_r.append(link_rPr)
                else:
                    # Ensure rPr comes first
                    link_r.remove(link_rPr)
                    link_r.insert(0, link_rPr)
                link_t = ET.SubElement(link_r, qn('a:t'))
                link_t.text = search_text

                # 3. Build AFTER run (regular continuation, no link)
                after_r = ET.Element(qn('a:r'))
                if orig_rPr is not None:
                    after_rPr = copy.deepcopy(orig_rPr)
                    for old in after_rPr.findall(qn('a:hlinkClick')):
                        after_rPr.remove(old)
                    after_r.append(after_rPr)
                after_t = ET.SubElement(after_r, qn('a:t'))
                after_t.text = after

                # 4. Insert link_r then after_r right after the original run
                orig_r.addnext(after_r)
                orig_r.addnext(link_r)
                return  # only first match
    # Not found — caller may want to know
    raise ValueError(f"search_text {search_text!r} not found in slide for hyperlinking")


PF_A, PF_ALPHA, PF_BETA = 3.155, 0.5, 0.3
PF_K_VALS = [100, 200, 300, 400]
# Table grid (TWELVE rows): extra L=250 step at the start surfaces the
# very steep early MPL (0→250 interval) without doubling the row count.
PF_L_VALS = [0, 250, 500, 1000, 1500, 2000, 2500, 3000, 3500, 4000, 4500, 5000]
# Chart grid (ELEVEN rows): uniform 500-step so the line-chart's
# categorical X-axis isn't visually distorted by the half-step at L=250.
# The L=250 point lives only in the TABLE; charts plot the smooth curve
# at uniform L intervals.
PF_L_VALS_CHART = [L for L in PF_L_VALS if L != 250]
# 2026-05-16: revised three times on the same day.
#   First pass:  β 0.5 → 0.4, A 0.5 → 1.25, L grid 0..10 000 step 1 000.
#   Second pass: β 0.4 → 0.3, A 1.25 → 3.155, L grid 0..5 000 step 500.
#   Third pass:  add intermittent L=250 row to table (per-worker MPL
#                strictly decreasing at every interval, including the
#                half-step from 0→250).  Q(400, 5 000) = 812.


def _pf_value(K, L):
    """Cobb-Douglas production function (integer-rounded cars per week)."""
    if K == 0 or L == 0:
        return 0
    return int(round(PF_A * K ** PF_ALPHA * L ** PF_BETA))


def _pf_table():
    """Full Q matrix indexed [row=L_index][col=K_index]."""
    return [[_pf_value(K, L) for K in PF_K_VALS] for L in PF_L_VALS]


def _add_compact_pf_table(slide, *, tbl_left, tbl_top, col_w_label=Inches(0.72),
                            col_w_data=Inches(0.55),
                            tbl_h=Inches(3.70),
                            font_size=11,
                            caption="Production-function table  (link)",
                            with_axes=True):
    """Insert the compact production-function table (same data as slide 10),
    with a drop-shadow rect behind and optional K/L axis labels + caption.

    Returns the table_shape so callers can position related elements.
    """
    Q_t = _pf_table()
    header_row = [""] + [str(K) for K in PF_K_VALS]
    rows_data = [header_row]
    for ri, L in enumerate(PF_L_VALS):
        rows_data.append([f"{L:,}"] + [str(v) for v in Q_t[ri]])

    rows = len(rows_data); cols = len(rows_data[0])
    tbl_w = col_w_label + col_w_data * 4

    _add_graphicframe_shadow(slide, tbl_left, tbl_top, tbl_w, tbl_h)
    tshape = slide.shapes.add_table(rows, cols, int(tbl_left), int(tbl_top),
                                      int(tbl_w), int(tbl_h))
    tbl = tshape.table
    tbl.columns[0].width = col_w_label
    for c in range(1, cols):
        tbl.columns[c].width = col_w_data

    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.01)
            cell.margin_bottom = Inches(0.01)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.name = "Calibri"
                    run.font.size = Pt(font_size)
                    if r == 0 or c == 0:
                        run.font.bold = True
                        run.font.color.rgb = WHITE
                    else:
                        run.font.color.rgb = NAVY
            if r == 0 or c == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE

    if caption:
        _add_text(slide,
                   tbl_left, tbl_top + tbl_h + Inches(0.08),
                   tbl_w, Inches(0.25),
                   caption,
                   size=11, italic=True, color=GRAY,
                   align=PP_ALIGN.CENTER, font="Calibri")

    if with_axes:
        _add_text(slide,
                   tbl_left + col_w_label, tbl_top - Inches(0.30),
                   col_w_data * 4, Inches(0.25),
                   "K  (robots)",
                   size=10, italic=True, color=NAVY,
                   align=PP_ALIGN.CENTER, font="Calibri")
        _add_text(slide,
                   tbl_left - Inches(0.75), tbl_top + tbl_h / 2 - Inches(0.15),
                   Inches(0.70), Inches(0.30),
                   "L  (workers)",
                   size=10, italic=True, color=NAVY,
                   align=PP_ALIGN.RIGHT,
                   anchor=MSO_ANCHOR.MIDDLE, font="Calibri")
    return tshape


# --------------------------------------------------------------------------
# Slide-10 user-added "Number of cars" callout group.
#
# This XML was hand-built in PowerPoint (oval circling a table cell,
# rectangle label "Number of / cars" on the right, slanted line with
# arrowhead from label to oval) and copied verbatim so that running
# `_build_clean_deck.py` reproduces it identically. Do NOT regenerate
# from python-pptx primitives – the styling (Whitney-Book font,
# stealth arrowhead, drop shadow on the second-line text) is hard to
# recreate via the python-pptx API and is preserved here as-is.
# --------------------------------------------------------------------------

GROUP_XML_SLIDE10 = '''<p:grpSp><p:nvGrpSpPr><p:cNvPr id="101" name="Group 16"><a:extLst><a:ext uri="{FF2B5EF4-FFF2-40B4-BE49-F238E27FC236}"><a16:creationId xmlns:a16="http://schemas.microsoft.com/office/drawing/2014/main" id="{CCFFE88A-C1F2-920E-75A7-DC37624177EC}"/></a:ext></a:extLst></p:cNvPr><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr><p:grpSpPr><a:xfrm><a:off x="5669321" y="3150554"/><a:ext cx="5054639" cy="1236280"/><a:chOff x="3902753" y="3429000"/><a:chExt cx="5054639" cy="1236280"/></a:xfrm></p:grpSpPr><p:sp><p:nvSpPr><p:cNvPr id="102" name="Oval 17"><a:extLst><a:ext uri="{FF2B5EF4-FFF2-40B4-BE49-F238E27FC236}"><a16:creationId xmlns:a16="http://schemas.microsoft.com/office/drawing/2014/main" id="{149BAB88-5F44-F462-B3EF-E90763362AA2}"/></a:ext></a:extLst></p:cNvPr><p:cNvSpPr/><p:nvPr/></p:nvSpPr><p:spPr><a:xfrm><a:off x="3902753" y="4279192"/><a:ext cx="969223" cy="386088"/></a:xfrm><a:prstGeom prst="ellipse"><a:avLst/></a:prstGeom><a:noFill/><a:ln w="28575"><a:solidFill><a:schemeClr val="tx1"/></a:solidFill></a:ln></p:spPr><p:style><a:lnRef idx="1"><a:schemeClr val="accent1"/></a:lnRef><a:fillRef idx="3"><a:schemeClr val="accent1"/></a:fillRef><a:effectRef idx="2"><a:schemeClr val="accent1"/></a:effectRef><a:fontRef idx="minor"><a:schemeClr val="lt1"/></a:fontRef></p:style><p:txBody><a:bodyPr rtlCol="0" anchor="ctr"/><a:lstStyle/><a:p><a:pPr algn="ctr"/><a:endParaRPr lang="en-US"/></a:p></p:txBody></p:sp><p:sp><p:nvSpPr><p:cNvPr id="103" name="Rectangle 8"><a:extLst><a:ext uri="{FF2B5EF4-FFF2-40B4-BE49-F238E27FC236}"><a16:creationId xmlns:a16="http://schemas.microsoft.com/office/drawing/2014/main" id="{64C69EA1-DAB2-A1C9-80EF-8CB4813BCBB1}"/></a:ext></a:extLst></p:cNvPr><p:cNvSpPr><a:spLocks noChangeArrowheads="1"/></p:cNvSpPr><p:nvPr/></p:nvSpPr><p:spPr bwMode="auto"><a:xfrm><a:off x="7532631" y="3429000"/><a:ext cx="1424761" cy="422371"/></a:xfrm><a:prstGeom prst="rect"><a:avLst/></a:prstGeom><a:noFill/><a:ln w="9525"><a:noFill/><a:miter lim="800000"/><a:headEnd/><a:tailEnd/></a:ln><a:effectLst/></p:spPr><p:txBody><a:bodyPr wrap="none" lIns="20048" tIns="28402" rIns="20048" bIns="28402"/><a:lstStyle/><a:p><a:pPr algn="ctr"><a:lnSpc><a:spcPts val="1700"/></a:lnSpc><a:tabLst><a:tab pos="481157" algn="l"/><a:tab pos="962315" algn="l"/><a:tab pos="1443472" algn="l"/></a:tabLst></a:pPr><a:r><a:rPr lang="en-US" b="0" dirty="0"><a:latin typeface="Whitney-Book" pitchFamily="50" charset="0"/><a:cs typeface="Whitney-Book" pitchFamily="50" charset="0"/></a:rPr><a:t>Number of</a:t></a:r></a:p><a:p><a:pPr algn="ctr"><a:lnSpc><a:spcPts val="1700"/></a:lnSpc><a:tabLst><a:tab pos="481157" algn="l"/><a:tab pos="962315" algn="l"/><a:tab pos="1443472" algn="l"/></a:tabLst></a:pPr><a:r><a:rPr lang="en-US" dirty="0"><a:effectLst><a:outerShdw blurRad="38100" dist="38100" dir="2700000" algn="tl"><a:srgbClr val="C0C0C0"/></a:outerShdw></a:effectLst><a:latin typeface="Whitney-Book" pitchFamily="50" charset="0"/><a:cs typeface="Whitney-Book" pitchFamily="50" charset="0"/></a:rPr><a:t>cars</a:t></a:r><a:endParaRPr lang="en-US" b="0" dirty="0"><a:effectLst><a:outerShdw blurRad="38100" dist="38100" dir="2700000" algn="tl"><a:srgbClr val="C0C0C0"/></a:outerShdw></a:effectLst><a:latin typeface="Whitney-Book" pitchFamily="50" charset="0"/><a:cs typeface="Whitney-Book" pitchFamily="50" charset="0"/></a:endParaRPr></a:p></p:txBody></p:sp><p:sp><p:nvSpPr><p:cNvPr id="104" name="Line 11"><a:extLst><a:ext uri="{FF2B5EF4-FFF2-40B4-BE49-F238E27FC236}"><a16:creationId xmlns:a16="http://schemas.microsoft.com/office/drawing/2014/main" id="{2D13A4FD-F41A-D960-86D4-A41F1454BC30}"/></a:ext></a:extLst></p:cNvPr><p:cNvSpPr><a:spLocks noChangeShapeType="1"/></p:cNvSpPr><p:nvPr/></p:nvSpPr><p:spPr bwMode="auto"><a:xfrm flipH="1"><a:off x="4871976" y="3727766"/><a:ext cx="2885453" cy="633968"/></a:xfrm><a:prstGeom prst="line"><a:avLst/></a:prstGeom><a:noFill/><a:ln w="34925"><a:solidFill><a:schemeClr val="tx1"/></a:solidFill><a:round/><a:headEnd type="none" w="sm" len="sm"/><a:tailEnd type="stealth" w="med" len="lg"/></a:ln><a:effectLst/></p:spPr><p:txBody><a:bodyPr wrap="none" lIns="96231" tIns="48116" rIns="96231" bIns="48116" anchor="ctr"/><a:lstStyle/><a:p><a:endParaRPr lang="en-US" sz="3200" b="0"><a:latin typeface="Whitney-Book" pitchFamily="50" charset="0"/><a:cs typeface="Whitney-Book" pitchFamily="50" charset="0"/></a:endParaRPr></a:p></p:txBody></p:sp></p:grpSp>'''


def _inject_raw_xml(slide, xml_str):
    """Append a raw XML element (e.g. a <p:grpSp>) to a slide's shape tree.

    The XML must be a single root element. We inject xmlns:p and xmlns:a
    declarations onto the root so it parses standalone — when re-serialised
    as part of the slide, lxml strips redundant declarations.
    """
    NS = (
        ' xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
    )
    # Insert namespaces right after the opening tag name
    tag_end = xml_str.find('>')
    space_in_tag = xml_str.find(' ', 0, tag_end)
    insert_at = space_in_tag if (space_in_tag != -1 and space_in_tag < tag_end) else tag_end
    wrapped = xml_str[:insert_at] + NS + xml_str[insert_at:]
    elem = ET.fromstring(wrapped)
    slide.shapes._spTree.append(elem)


def slide_10(prs):
    """Rivian's production function: weekly output of the R1 line.

    A 2D table showing weekly output as a function of (workers, robots).
    Values come from Cobb-Douglas Q = 4 · K^0.3 · L^0.5 – chosen so that
    BOTH MPL (across L for fixed K) AND MPK (across K for fixed L) are
    STRICTLY diminishing. (The original-deck table had constant MPK across
    the first three K-steps, which technically violated diminishing
    returns to K.)
    """
    def draw(slide):
        Q = _pf_table()
        # Build display table: header row + data rows. Each cell renders
        # workforce sizes with thousands separator for legibility.
        header = [""] + [str(K) for K in PF_K_VALS]
        rows_data = [header]
        for ri, L in enumerate(PF_L_VALS):
            rows_data.append([f"{L:,}"] + [str(v) for v in Q[ri]])

        rows = len(rows_data)
        cols = len(rows_data[0])

        # Columns sized to leave ~0.5 cm padding either side of the longest
        # number at 14pt Calibri ("10,000" in col 0; "1000" in data cols).
        col_w_label = Inches(1.00)
        col_w_data = Inches(0.80)
        data_cols_w = col_w_data * 4
        tbl_w = col_w_label + data_cols_w
        tbl_h = Inches(4.00)           # ~0.33" per row × 12 rows
        tbl_top = Inches(2.20)         # shifted up from 2.45 on 2026-05-15
                                        # to make room for the Concept-
                                        # explanation callout at the bottom
        tbl_left = int((SLIDE_W - tbl_w) / 2)   # centre horizontally

        # Soft drop shadow rectangle BEHIND the table (graphicFrames can't host shadow).
        _add_graphicframe_shadow(slide, tbl_left, tbl_top, tbl_w, tbl_h)

        table_shape = slide.shapes.add_table(rows, cols, tbl_left, tbl_top,
                                              tbl_w, tbl_h)
        tbl = table_shape.table
        tbl.columns[0].width = col_w_label
        for c in range(1, cols):
            tbl.columns[c].width = col_w_data

        cell_pad = Inches(0.20)        # ≈ 0.5 cm horizontal
        cell_pad_v = Inches(0.02)
        for r, row in enumerate(rows_data):
            for c, val in enumerate(row):
                cell = tbl.cell(r, c)
                cell.margin_left = cell_pad
                cell.margin_right = cell_pad
                cell.margin_top = cell_pad_v
                cell.margin_bottom = cell_pad_v
                cell.text = str(val)
                tf = cell.text_frame
                for p in tf.paragraphs:
                    p.alignment = PP_ALIGN.CENTER
                    for run in p.runs:
                        run.font.name = "Calibri"
                        run.font.size = Pt(14)
                        if r == 0 or c == 0:
                            run.font.bold = True
                            run.font.color.rgb = WHITE
                        else:
                            run.font.color.rgb = NAVY
                if r == 0 or c == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = NAVY
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = WHITE

        # --- Direction-of-increase arrows (gold block arrows) ---
        # Right-arrow above the data columns – more robots →
        # 2026-05-17: arrow y nudged 1.80 → 1.860 per manual edit.
        data_cols_left = tbl_left + col_w_label
        top_arrow_h = Inches(0.30)
        top_arrow_top = Inches(1.860)
        top_arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            int(data_cols_left), int(top_arrow_top),
            int(data_cols_w), int(top_arrow_h),
        )
        top_arrow.fill.solid()
        top_arrow.fill.fore_color.rgb = GOLD
        top_arrow.line.fill.background()
        top_arrow.shadow.inherit = False

        # Down-arrow to the left of the data rows – more workers ↓
        # 2026-05-17: y and height hand-tweaked.
        data_rows_top = Inches(2.508)
        data_rows_h = Inches(3.692)
        left_arrow_w = Inches(0.30)
        left_arrow_left = tbl_left - left_arrow_w - Inches(0.12)
        left_arrow = slide.shapes.add_shape(
            MSO_SHAPE.DOWN_ARROW,
            int(left_arrow_left), int(data_rows_top),
            int(left_arrow_w), int(data_rows_h),
        )
        left_arrow.fill.solid()
        left_arrow.fill.fore_color.rgb = GOLD
        left_arrow.line.fill.background()
        left_arrow.shadow.inherit = False

        # --- Axis labels (above top arrow / left of down arrow) ---
        # 2026-05-17: K label moved down to overlay the arrow header per
        # manual edit (top 1.30 → 1.610).
        top_label_h = Inches(0.40)
        top_label_y = Inches(1.610)
        _add_text(slide, int(data_cols_left), int(top_label_y),
                   int(data_cols_w), int(top_label_h),
                   "Number of robots (K)",
                   size=18, bold=True, color=NAVY,
                   align=PP_ALIGN.CENTER, font="Calibri")
        label_left_w = Inches(2.0)
        label_left_h = Inches(0.8)
        label_left_x = int(left_arrow_left - label_left_w - Inches(0.10))
        # Hand-nudge: centre with a -0.10" vertical offset to match user edit.
        label_left_y = int(data_rows_top + data_rows_h / 2 - label_left_h / 2
                            - Inches(0.10))
        _add_text(slide, label_left_x, label_left_y,
                   int(label_left_w), int(label_left_h),
                   "Number of\nworkers (L)",
                   size=18, bold=True, color=NAVY,
                   align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE,
                   font="Calibri")
        # Concept-explanation callout below the table — cream-fill rounded
        # rect.  Two short lines, 17 pt navy bold.
        # 2026-05-17: narrowed and shifted right to clear the "Number of
        # cars" annotation group (was width 11.0 centred; now 7.086 at
        # x=3.359).
        cap_w = Inches(7.086)
        cap_h = Inches(0.665)
        cap_x = Inches(3.359)
        cap_y = Inches(6.320)
        _add_convention_box(
            slide, cap_x, cap_y, cap_w, cap_h,
            runs=[
                ("Output = cars per week",
                 {'size': 17, 'bold': True, 'color': NAVY}),
                ("MPL falls down each column;  MPK falls along each row",
                 {'size': 17, 'bold': True, 'color': NAVY, 'newline': True}),
            ],
            size=17, align=PP_ALIGN.CENTER,
        )
        # --- Inject the user-added "Number of cars" callout group ---
        _inject_raw_xml(slide, GROUP_XML_SLIDE10)

    s = make_diagram_slide(
        prs, page_num=11,
        section_tag=SECTION_TAG_P1,
        title="Rivian's Production Function:  R1 Line Weekly Output",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "Rivian's R1 line weekly output as a function of workers (L) and "
        "robots (K).  Built from the Cobb-Douglas Q = 3.155 · √K · L^0.3 – "
        "strictly diminishing in each input individually.  Look down a "
        "column: MPL falls as you add labor with capital fixed.  Look "
        "along a row: MPK falls as you add robots with labor fixed.  "
        "Strict diminishing returns in both directions – the textbook "
        "story.  Per-worker MPL is also strictly diminishing across the "
        "non-uniform L grid (extra L=250 step at the start)."
    ))


def slide_11(prs):
    """Plot total output Q vs. L for four plant-capital levels.

    Native line-with-markers chart driven by the same Cobb-Douglas table
    on slide 10. Four K series (100 / 200 / 300 / 400 robots) – matches
    the original deck's chart layout, including marker shapes per series
    and the in-plot legend in the top-left corner.
    """
    SER_COLORS = [
        RGBColor(0x2E, 0x75, 0xB6),  # blue   – K=100
        RGBColor(0xC0, 0x50, 0x4D),  # red    – K=200
        RGBColor(0x80, 0x80, 0x80),  # gray   – K=300
        RGBColor(0xE6, 0xB8, 0x00),  # gold   – K=400
    ]
    SER_MARKERS = ['circle', 'triangle', 'square', 'diamond']

    def draw(slide):
        chart_data = CategoryChartData()
        # 2026-05-16 (third pass): use PF_L_VALS (with L=250) — per user
        # request, all illustrations carry the extra half-step at the
        # start, even though the categorical X-axis spaces it equally with
        # the other intervals.
        chart_data.categories = [f"{L:,}" for L in PF_L_VALS]
        for K in PF_K_VALS:
            series_vals = [_pf_value(K, L) for L in PF_L_VALS]
            chart_data.add_series(f"K = {K}", series_vals)

        # Chart frame: 8.4" wide, top y=1.30 (just below the title divider).
        # Height trimmed to 4.80" on 2026-05-15 to make room for the new
        # cream Concept-explanation callout below (which is larger than the
        # earlier thin banner).  Bottom at y≈6.10.
        chart_w = Inches(8.4)
        chart_h = Inches(4.80)
        chart_top = Inches(1.30)
        chart_left = Inches(2.636)
        # Drop-shadow rectangle behind the chart.
        _add_graphicframe_shadow(slide, chart_left, chart_top, chart_w, chart_h)
        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE,  # markers added per series below
            chart_left, chart_top, chart_w, chart_h,
            chart_data,
        )
        chart = chart_shape.chart

        # No internal "Output per Week" title – redundant with the slide
        # action title; dropping it reclaims ~0.4" of plot-area height.
        chart.has_title = False

        # --- Tighten the white margin around the plot ---------------------
        # By default PowerPoint reserves a generous border between the chart
        # frame and the plot region.  Force a manualLayout (layoutTarget=
        # "inner") so the inner plot fills ~88 % × 82 % of the chart frame,
        # leaving room only for the y-axis labels (left), x-axis labels +
        # title (bottom), and a thin top margin.  Added 2026-05-15.
        # chart._element is the <c:chartSpace>; plotArea lives at
        # chartSpace/chart/plotArea — navigate two levels down.
        chart_el = chart._element.find(qn('c:chart'))
        plot_el = chart_el.find(qn('c:plotArea')) if chart_el is not None else None
        if plot_el is not None:
            for old in plot_el.findall(qn('c:layout')):
                plot_el.remove(old)
            pl_layout = ET.Element(qn('c:layout'))
            pl_ml = ET.SubElement(pl_layout, qn('c:manualLayout'))
            ltgt = ET.SubElement(pl_ml, qn('c:layoutTarget')); ltgt.set('val', 'inner')
            xM = ET.SubElement(pl_ml, qn('c:xMode')); xM.set('val', 'edge')
            yM = ET.SubElement(pl_ml, qn('c:yMode')); yM.set('val', 'edge')
            xv = ET.SubElement(pl_ml, qn('c:x')); xv.set('val', '0.10')
            yv = ET.SubElement(pl_ml, qn('c:y')); yv.set('val', '0.03')
            wv = ET.SubElement(pl_ml, qn('c:w')); wv.set('val', '0.88')
            hv = ET.SubElement(pl_ml, qn('c:h')); hv.set('val', '0.82')
            plot_el.insert(0, pl_layout)

        # Native legend, positioned inside the plot area (top-left).
        # Font bumped 12 → 13 pt and box dims bumped on 2026-05-15 per
        # user request — agenda reads a touch larger.
        chart.has_legend = True
        chart.legend.font.name = "Calibri"
        chart.legend.font.size = Pt(13)
        chart.legend.font.color.rgb = NAVY
        chart.legend.include_in_layout = False
        # Force legend to a manual layout inside the plot area (top-left).
        # python-pptx doesn't expose this, so write the layout XML directly.
        leg_el = chart.legend._element
        for old in leg_el.findall(qn('c:layout')):
            leg_el.remove(old)
        for old in leg_el.findall(qn('c:legendPos')):
            leg_el.remove(old)
        # Insert legendPos = 'tr' (anything works, layout below overrides)
        pos_el = ET.SubElement(leg_el, qn('c:legendPos'))
        pos_el.set('val', 'tr')
        leg_el.remove(pos_el)
        leg_el.insert(0, pos_el)
        # Insert <c:layout><c:manualLayout>… positioning legend in upper-left.
        # x bumped from 0.08 → 0.18 on 2026-05-15 so the legend clears the
        # y-axis area; h shrunk from 0.32 to 0.18 to keep the white-fill box
        # tight around the four legend entries (4 lines × ~12pt at 70 % LS).
        layout = ET.SubElement(leg_el, qn('c:layout'))
        ml = ET.SubElement(layout, qn('c:manualLayout'))
        xMode = ET.SubElement(ml, qn('c:xMode')); xMode.set('val', 'edge')
        yMode = ET.SubElement(ml, qn('c:yMode')); yMode.set('val', 'edge')
        x_el = ET.SubElement(ml, qn('c:x')); x_el.set('val', '0.18')
        y_el = ET.SubElement(ml, qn('c:y')); y_el.set('val', '0.05')
        w_el = ET.SubElement(ml, qn('c:w')); w_el.set('val', '0.17')
        h_el = ET.SubElement(ml, qn('c:h')); h_el.set('val', '0.24')
        # Re-order: legendPos must precede layout (already done by insert(0)).
        # Move <c:layout> right after <c:legendPos>.
        leg_el.remove(layout)
        leg_el.insert(list(leg_el).index(pos_el) + 1, layout)

        # --- Solid white fill + thin navy border on the legend box ---
        # Punches out the dashed gridlines underneath so the four series
        # labels read cleanly. Schema order: legendPos, legendEntry, layout,
        # overlay, spPr, txPr.  We ensure spPr sits after layout and before
        # any txPr that python-pptx may have created.
        for old in leg_el.findall(qn('c:spPr')):
            leg_el.remove(old)
        leg_spPr = ET.Element(qn('c:spPr'))
        sp_fill = ET.SubElement(leg_spPr, qn('a:solidFill'))
        sp_clr = ET.SubElement(sp_fill, qn('a:srgbClr'))
        sp_clr.set('val', 'FFFFFF')
        sp_ln = ET.SubElement(leg_spPr, qn('a:ln'))
        sp_ln.set('w', '6350')                 # 0.5 pt
        sp_lf = ET.SubElement(sp_ln, qn('a:solidFill'))
        sp_lc = ET.SubElement(sp_lf, qn('a:srgbClr'))
        sp_lc.set('val', '0B2B4E')             # NAVY
        # Insert immediately after c:layout (and any c:overlay that may exist)
        anchor = layout
        ovr = leg_el.find(qn('c:overlay'))
        if ovr is not None and list(leg_el).index(ovr) > list(leg_el).index(layout):
            anchor = ovr
        anchor.addnext(leg_spPr)

        # --- Tighter line spacing between entries (70%) ---
        # python-pptx already creates a c:txPr when font properties are set
        # above. We poke a:lnSpc into its first a:pPr; if no txPr exists,
        # build a minimal one.
        txPr = leg_el.find(qn('c:txPr'))
        if txPr is None:
            txPr = ET.Element(qn('c:txPr'))
            ET.SubElement(txPr, qn('a:bodyPr'))
            ET.SubElement(txPr, qn('a:lstStyle'))
            p_el = ET.SubElement(txPr, qn('a:p'))
            ET.SubElement(p_el, qn('a:pPr'))
            ET.SubElement(p_el, qn('a:endParaRPr'))
            leg_spPr.addnext(txPr)
        else:
            # Ensure txPr is after spPr (schema)
            leg_el.remove(txPr)
            leg_spPr.addnext(txPr)
        p_el = txPr.find(qn('a:p'))
        if p_el is None:
            p_el = ET.SubElement(txPr, qn('a:p'))
        pPr_el = p_el.find(qn('a:pPr'))
        if pPr_el is None:
            pPr_el = ET.Element(qn('a:pPr'))
            p_el.insert(0, pPr_el)
        for old in pPr_el.findall(qn('a:lnSpc')):
            pPr_el.remove(old)
        lnSpc = ET.Element(qn('a:lnSpc'))
        spcPct = ET.SubElement(lnSpc, qn('a:spcPct'))
        spcPct.set('val', '70000')             # 70%
        pPr_el.insert(0, lnSpc)

        # Axes
        cat = chart.category_axis
        cat.tick_labels.font.name = "Calibri"
        cat.tick_labels.font.size = Pt(11)
        cat.tick_labels.font.color.rgb = NAVY
        cat.has_title = True
        cat.axis_title.text_frame.text = "Number of Workers"
        ar = cat.axis_title.text_frame.paragraphs[0].runs[0]
        ar.font.name = "Calibri"; ar.font.size = Pt(14)
        ar.font.bold = True; ar.font.italic = True; ar.font.color.rgb = NAVY

        # Light-grey dashed vertical gridlines at each category tick
        # (1,000 / 2,000 / … / 10,000 workers). Schema order: majorGridlines
        # sits between c:axPos and c:title, so insert right after c:axPos.
        cat_el = cat._element
        for old in cat_el.findall(qn('c:majorGridlines')):
            cat_el.remove(old)
        gridlines = ET.Element(qn('c:majorGridlines'))
        sp = ET.SubElement(gridlines, qn('c:spPr'))
        ln = ET.SubElement(sp, qn('a:ln'))
        ln.set('w', '9525')                        # 0.75 pt
        ln.set('cap', 'flat'); ln.set('cmpd', 'sng'); ln.set('algn', 'ctr')
        fill = ET.SubElement(ln, qn('a:solidFill'))
        clr = ET.SubElement(fill, qn('a:srgbClr'))
        clr.set('val', 'C8CDD3')                   # RULE light grey
        dash = ET.SubElement(ln, qn('a:prstDash'))
        dash.set('val', 'dash')
        axpos = cat_el.find(qn('c:axPos'))
        axpos.addnext(gridlines)

        val = chart.value_axis
        val.tick_labels.font.name = "Calibri"
        val.tick_labels.font.size = Pt(11)
        val.tick_labels.font.color.rgb = NAVY
        val.minimum_scale = 0
        val.maximum_scale = 800
        val.major_unit = 100
        val.has_title = True
        val.axis_title.text_frame.text = "Cars per Week"
        ar = val.axis_title.text_frame.paragraphs[0].runs[0]
        ar.font.name = "Calibri"; ar.font.size = Pt(14)
        ar.font.bold = True; ar.font.italic = True; ar.font.color.rgb = NAVY

        # Align X-axis category labels (and data points) with the tick
        # marks rather than placing them in the gaps between ticks.
        _align_x_labels_with_ticks(val)

        # Light-grey dashed horizontal gridlines at each value tick
        # (100, 200, … 1000 cars per week).
        val_el = val._element
        for old in val_el.findall(qn('c:majorGridlines')):
            val_el.remove(old)
        v_gl = ET.Element(qn('c:majorGridlines'))
        sp = ET.SubElement(v_gl, qn('c:spPr'))
        ln = ET.SubElement(sp, qn('a:ln'))
        ln.set('w', '9525')
        ln.set('cap', 'flat'); ln.set('cmpd', 'sng'); ln.set('algn', 'ctr')
        fill = ET.SubElement(ln, qn('a:solidFill'))
        clr = ET.SubElement(fill, qn('a:srgbClr')); clr.set('val', 'C8CDD3')
        dash = ET.SubElement(ln, qn('a:prstDash')); dash.set('val', 'dash')
        v_axpos = val_el.find(qn('c:axPos'))
        v_axpos.addnext(v_gl)

        # Style each series: distinct color + marker shape.
        for idx, series in enumerate(chart.series):
            line = series.format.line
            line.color.rgb = SER_COLORS[idx]
            line.width = Pt(2.0)
            ser_xml = series._element
            # Marker block: <c:marker><c:symbol val="…"/><c:size val="7"/>
            #               <c:spPr>(solid fill + outline)</c:spPr></c:marker>
            for old in ser_xml.findall(qn('c:marker')):
                ser_xml.remove(old)
            marker = ET.SubElement(ser_xml, qn('c:marker'))
            sym = ET.SubElement(marker, qn('c:symbol'))
            sym.set('val', SER_MARKERS[idx])
            size_el = ET.SubElement(marker, qn('c:size'))
            size_el.set('val', '7')
            spPr = ET.SubElement(marker, qn('c:spPr'))
            fill = ET.SubElement(spPr, qn('a:solidFill'))
            rgb = ET.SubElement(fill, qn('a:srgbClr'))
            r, g, b = SER_COLORS[idx][0], SER_COLORS[idx][1], SER_COLORS[idx][2]
            rgb.set('val', f'{r:02X}{g:02X}{b:02X}')
            ln = ET.SubElement(spPr, qn('a:ln'))
            ln_fill = ET.SubElement(ln, qn('a:solidFill'))
            ln_rgb = ET.SubElement(ln_fill, qn('a:srgbClr'))
            ln_rgb.set('val', f'{r:02X}{g:02X}{b:02X}')
            # Move marker block before c:smooth / after c:spPr ordering
            # (schema: order, idx, tx, spPr, marker, …)
            # python-pptx adds elements in normal sub-element order;
            # if needed, reorder explicitly:
            order = ['c:idx', 'c:order', 'c:tx', 'c:spPr', 'c:marker',
                     'c:dPt', 'c:dLbls', 'c:trendline', 'c:errBars',
                     'c:cat', 'c:val', 'c:smooth']
            children = list(ser_xml)
            children.sort(key=lambda el: order.index(el.tag.replace(
                '{http://schemas.openxmlformats.org/drawingml/2006/chart}',
                'c:')) if el.tag.replace(
                '{http://schemas.openxmlformats.org/drawingml/2006/chart}',
                'c:') in order else 999)
            for c in children:
                ser_xml.remove(c)
            for c in children:
                ser_xml.append(c)
            # Disable smoothing so curves are straight segments between points
            # (matches the original chart's piecewise-linear look).
            for sm in ser_xml.findall(qn('c:smooth')):
                ser_xml.remove(sm)
            smooth = ET.SubElement(ser_xml, qn('c:smooth'))
            smooth.set('val', '0')

        # Concept-explanation callout below the chart — cream-fill rounded
        # rect.  2026-05-17: width + x aligned with the chart's left edge
        # and width so the callout reads as the chart's footer band.
        banner_w = Inches(8.286)
        banner_h = Inches(0.780)
        banner_x = Inches(2.636)
        banner_y = Inches(6.235)
        _add_convention_box(
            slide, banner_x, banner_y, banner_w, banner_h,
            runs=[
                ("Each curve flattens as L rises  (diminishing MPL)",
                 {'size': 17, 'bold': True, 'color': NAVY}),
                ("The vertical distance between curves narrows as K rises  (diminishing MPK)",
                 {'size': 17, 'bold': True, 'color': NAVY, 'newline': True}),
            ],
            size=17, align=PP_ALIGN.CENTER,
            pad_h=Inches(0.15), pad_v=Inches(0.04),
        )

    s = make_diagram_slide(
        prs, page_num=12,
        section_tag=SECTION_TAG_P1,
        title="Plotting Total Output:  Q vs. L on the R1 Line",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "Same numbers as the table on the previous slide, now plotted. "
        "Four lines for four plant-capital levels (100, 200, 300, 400 "
        "robots). Each line rises but flattens – diminishing MPL. The "
        "vertical spacing between adjacent lines narrows as K grows – "
        "diminishing MPK. Memorise this shape: it shows up in every "
        "production context, from factories to consulting teams."
    ))


def slide_short_run_agenda(prs):
    """Section divider – Part 1.b: Short Run (page 13).

    2026-05-19: NEW slide inserted between the production-function
    motivation (pages 9-12) and the short-run hiring mechanics
    (page 14 onwards).  User signal that we've crossed from "what is
    a production function" into "how to use it in the short run".

    Same Layout 2 / agenda view as slide_7 (Part 1 navy, Part 2 faded)
    but with the b sub of Part 1 as the only fully-navy bullet — subs
    a (Production Function) and c (Long Run) render in faded gray to
    make the "you are here" cue unambiguous.
    """
    s = make_section_agenda(
        prs, page_num=13,
        current_part_idx=0,
        current_sub_idx=1,
        section_tag=SECTION_TAG_DIV,
        title="Part 1.b:  Short Run – Hiring Decisions",
    )
    _set_notes(s, (
        "Quick orientation. We've established what a production function "
        "is and how it behaves. Now we shift into the short run – capital "
        "is fixed, labor is the lever – and ask the operating question: "
        "how many workers should the firm hire? The next several slides "
        "build the marginal-product machinery we'll use to answer that."
    ))


def slide_12(prs):
    """Short Run: Marginal Product of Labor — concept intro.

    Mirrors the source deck's original slide 16: re-anchor the short-run
    framing (K fixed, L flexible), introduce the MPL concept name in
    accent blue, give the formal change/change definition with "change"
    flagged in accent red italic, and close with the canonical
    ΔQ/ΔL formula in big OMML below.
    """
    ACCENT_BLUE = RGBColor(0x00, 0x70, 0xC0)
    DARK_YELLOW = RGBColor(0xB8, 0x86, 0x0B)   # was red; per user request

    def _styled_run(p, text, *, size=24, bold=False, italic=False,
                    color=NAVY, font="Calibri"):
        r = p.add_run()
        r.text = text
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        return r

    def _block(slide, left, top, width, height, runs, *,
               align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(int(left), int(top),
                                        int(width), int(height))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
        tf.margin_top = Inches(0.0);  tf.margin_bottom = Inches(0.0)
        p = tf.paragraphs[0]
        p.alignment = align
        for text, opts in runs:
            _styled_run(p, text, **opts)
        return tb

    def draw(slide):
        indent = Inches(0.45)
        # 1) MAIN BULLET — "Typically in the short run:"  (28pt bold navy,
        #    prefixed with a navy ▪ glyph)
        _block(slide, MARGIN, Inches(1.85), RULE_W, Inches(0.45), [
            ("▪  ",                {'size': 28, 'bold': True}),
            ("Typically",          {'size': 28, 'bold': True}),
            (" in the short run:", {'size': 28, 'bold': True}),
        ])
        # 2) SUB-BULLETS — Capital is fixed (K̅) / Labor is flexible (L)
        #    Two paragraphs with inline OMML for the symbols.  Sub-text
        #    bumped 24 → 26 pt on 2026-05-15 (sub-bullets were too small).
        _add_mixed_textbox(slide,
                            MARGIN + indent, Inches(2.45),
                            RULE_W - indent, Inches(1.20),
                            [
                                ('text', "–  ", {'size': 26}),
                                ('text', "Capital is ", {'size': 26}),
                                ('text', "fixed", {'size': 26, 'bold': True}),
                                ('text', "  (", {'size': 26}),
                                ('omml', _omml_acc_overline('K'), {'size': 26}),
                                ('text', ")", {'size': 26}),
                                ('break', '', {}),
                                ('text', "–  ", {'size': 26}),
                                ('text', "Labor is flexible  (", {'size': 26}),
                                ('omml', _omml_run('L'), {'size': 26}),
                                ('text', ")", {'size': 26}),
                            ],
                            default_size=26, default_color=NAVY)

        # 3) MAIN BULLET — "Important Concept:  Marginal Product of Labor"
        _block(slide, MARGIN, Inches(4.00), RULE_W, Inches(0.45), [
            ("▪  ",                        {'size': 28, 'bold': True}),
            ("Important Concept:  ",       {'size': 28, 'bold': True}),
            ("Marginal Product of Labor",  {'size': 28, 'bold': True,
                                             'color': ACCENT_BLUE}),
        ])
        # 4) SUB-BULLET — formal definition with "change" emphasised,
        #    26pt (bumped from 24 — see same note as above), indented.
        def_tb = slide.shapes.add_textbox(
            int(MARGIN + indent), int(Inches(4.60)),
            int(RULE_W - indent), int(Inches(0.75)))
        def_tf = def_tb.text_frame
        def_tf.word_wrap = True
        def_tf.margin_left = Inches(0.05); def_tf.margin_right = Inches(0.05)
        def_tf.margin_top = Inches(0); def_tf.margin_bottom = Inches(0)
        p = def_tf.paragraphs[0]
        _styled_run(p, "–  ", size=26)
        _styled_run(p, "The ", size=26)
        _styled_run(p, "marginal product of labor ", size=26,
                    color=ACCENT_BLUE)
        _styled_run(p, "is the ", size=26)
        _styled_run(p, "change", size=26, italic=True, color=DARK_YELLOW)
        _styled_run(p, " in output due to a ", size=26)
        _styled_run(p, "change", size=26, italic=True, color=DARK_YELLOW)
        _styled_run(p, " in labor input:", size=26)

        # 5) Big OMML formula MPL = ΔQ / ΔL  (36pt, blue)
        mpl     = _omml_sub(_omml_run('MP'), _omml_run('L'))
        delta_q = _omml_text('Δ') + _omml_run('Q')
        delta_l = _omml_text('Δ') + _omml_run('L')
        frac    = _omml_frac(delta_q, delta_l)
        omml_full = mpl + _omml_text(' = ') + frac
        _add_math_equation(slide,
                            left=Inches(4.7), top=Inches(5.55),
                            width=Inches(4.0), height=Inches(1.25),
                            omml_content=omml_full,
                            size_pt=36, color=ACCENT_BLUE)

    s = make_diagram_slide(
        prs, page_num=14,
        section_tag=SECTION_TAG_P1,
        title="Short Run:  Marginal Product of Labor",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "Re-anchor the short-run framing: capital is fixed (we write the "
        "bar over K to emphasise that K is held constant) while labor is "
        "flexible. Now introduce the marginal product of labor formally: "
        "MPL is the CHANGE in output due to a CHANGE in labor input. The "
        "formula MPL = ΔQ / ΔL is shorthand for \"how much extra output "
        "do I get from one more worker?\" – exactly the question every "
        "plant manager asks every shift.  Next slide: the actual numbers "
        "from the Rivian production function show MPL declining."
    ))


def slide_mpl_data(prs):
    """Numerical MPL example — matches original slide 17 of the source deck.

    Fixes K = 100 robots, walks L through {0, 500, 1k, 2k, 3k}, and
    builds the L | K | Q | ΔL | ΔQ | MPL table from the same Cobb-Douglas
    function as slides 10 and 11.  MPL falls strictly down the column,
    making the "diminishing marginal product" lesson visible numerically.
    A right-hand Convention box explains the step-by-step ΔL / ΔQ rule.
    """
    ACCENT_BLUE = RGBColor(0x00, 0x70, 0xC0)
    ACCENT_RED  = RGBColor(0xFF, 0x00, 0x00)
    MPL_FILL    = RGBColor(0xFF, 0xF5, 0xE0)   # cream highlight for MPL column
    CONV_FILL   = RGBColor(0xFD, 0xF6, 0xE6)   # softer cream for callout

    # Per-column body-cell colors (header row stays white-on-navy).
    # Column order on 2026-05-15: L | K | Q | ΔQ | ΔL | MPL  (ΔQ before
    # ΔL, swapped from the original ordering per user request).
    BLACK_NUM = RGBColor(0x00, 0x00, 0x00)
    RED_NUM   = RGBColor(0xC0, 0x00, 0x00)
    GREEN_NUM = RGBColor(0x1B, 0x5E, 0x20)        # darker / deeper green
    BLUE_NUM  = ACCENT_BLUE
    COL_COLORS = [BLACK_NUM, RED_NUM, BLACK_NUM, GREEN_NUM, GREEN_NUM, BLUE_NUM]

    K_FIX = 100
    # 2026-05-18: extended L_GRID to match slide 11's worker steps up to
    # L = 2,500 (was [0, 250, 500, 1000, 2000, 3000] — top of grid moved
    # from 3,000 to 2,500 and intermediate stops 1,500 / 2,000 added so
    # every step is a clean 500 after the initial 250-worker steps).
    L_GRID = [0, 250, 500, 1000, 1500, 2000, 2500]

    # 2026-05-18 (manual): per-interval Y centres for ΔQ / ΔL / MPL
    # floats and their accompanying down-arrows and wavy connectors.
    # The user dragged each row's float spacing by hand in PowerPoint
    # so the floats sit visually at the midpoints between adjacent Q
    # cells.  Raw row XML still reports 0.3375"/row, but PowerPoint
    # renders the (resized) 2.8575" table by stretching rows — these
    # Y values were sampled directly from the canonical deck.
    FLOAT_CENTER_Y = [
        None,             # i=0 placeholder (no float between header and L=0 row)
        Inches(3.525),    # i=1: between L=0 and L=250    (= mathematical boundary)
        Inches(3.893),    # i=2: between L=250 and L=500  (+0.031" vs grid)
        Inches(4.250),    # i=3: between L=500 and L=1000 (+0.050" vs grid)
        Inches(4.598),    # i=4: between L=1000 and L=1500 (+0.061" vs grid)
        Inches(4.965),    # i=5: between L=1500 and L=2000 (+0.090" vs grid)
        Inches(5.332),    # i=6: between L=2000 and L=2500 (+0.119" vs grid)
    ]

    def draw(slide):
        # Main bullet — replaces the old centred italic captions with a
        # proper bullet structure (per user request 2026-05-15).
        _add_mixed_textbox(slide,
                            MARGIN, Inches(1.85),
                            RULE_W, Inches(0.45),
                            [
                                ('text', "▪  ",
                                 {'size': 24, 'bold': True, 'color': NAVY}),
                                ('text', "Example:  MPL from Rivian Production function.",
                                 {'size': 24, 'bold': True, 'color': NAVY}),
                            ],
                            align=PP_ALIGN.LEFT,
                            default_size=24, default_color=NAVY)

        # Sub-bullet with inline OMML K̅
        _add_mixed_textbox(slide,
                            MARGIN + Inches(0.45), Inches(2.35),
                            RULE_W - Inches(0.45), Inches(0.40),
                            [
                                ('text', "–  ", {'size': 22, 'color': NAVY}),
                                ('text', "Fix capital at  ",
                                 {'size': 22, 'color': NAVY}),
                                ('omml', _omml_acc_overline('K'),
                                 {'size': 22}),
                                ('text', "  =  100",
                                 {'size': 22, 'bold': True, 'color': NAVY}),
                            ],
                            align=PP_ALIGN.LEFT,
                            default_size=22, default_color=NAVY)

        # ---- Table (6 columns, including the new K column) ----
        # The Δ-columns (ΔL, ΔQ, MPL) are rendered specially on
        # 2026-05-15: their cells are blank inside the table, and the
        # values are drawn as floating textboxes positioned at the
        # BOUNDARY between two adjacent rows.  This visually illustrates
        # the convention that each Δ is computed relative to the
        # previous (initial) point — values live "between" rows, not on
        # them.  MPL floats also get the cream MPL_FILL background;
        # ΔL/ΔQ floats are transparent (green numbers on white).
        Q = [_pf_value(K_FIX, L) for L in L_GRID]
        dL_values  = [None]
        dQ_values  = [None]
        mpl_values = [None]
        rows_data = [["L", "K", "Q", "ΔQ", "ΔL", "MPL"]]
        for i, L in enumerate(L_GRID):
            row = [f"{L:,}", f"{K_FIX}", f"{Q[i]:,}"]
            if i == 0:
                row += ["", "", ""]                       # all 3 Δ cells empty
            else:
                dL = L_GRID[i] - L_GRID[i-1]
                dQ = Q[i] - Q[i-1]
                mpl = dQ / dL
                row += ["", "", ""]                       # all 3 Δ cells empty
                dL_values.append(f"{dL:,}")
                dQ_values.append(f"{dQ}")
                mpl_values.append(f"{mpl:.3f}")
            rows_data.append(row)

        rows = len(rows_data); cols = len(rows_data[0])
        col_widths = [Inches(0.80), Inches(0.65),
                       Inches(0.80), Inches(0.85),
                       Inches(0.75), Inches(0.95)]
        tbl_w = sum(col_widths)
        # 2026-05-18: bumped tbl_h from 2.55" → 2.70" to accommodate one
        # extra L-row (now 8 rows incl. header) while keeping row_h close
        # to the original (~0.338" vs. old ~0.364").
        # 2026-05-18 (later, manual): user resized the table in PowerPoint
        # to 2.8575" (rows still 0.3375" in XML, ~0.358" rendered).
        # Matching that here so floats positioned by FLOAT_CENTER_Y sit
        # on the rendered row boundaries.
        tbl_h = Inches(2.8575)
        tbl_top = Inches(2.85)
        # Table no longer centred – keep it on the LEFT so a Convention
        # callout fits to its right.
        tbl_left = Inches(0.80)
        _add_graphicframe_shadow(slide, tbl_left, tbl_top, tbl_w, tbl_h)
        tshape = slide.shapes.add_table(rows, cols, tbl_left, tbl_top,
                                          tbl_w, tbl_h)
        tbl = tshape.table
        for ci, w in enumerate(col_widths):
            tbl.columns[ci].width = w

        cell_pad_h = Inches(0.10)
        for r, row in enumerate(rows_data):
            for c, val in enumerate(row):
                cell = tbl.cell(r, c)
                cell.margin_left = cell_pad_h
                cell.margin_right = cell_pad_h
                cell.margin_top = Inches(0.03)
                cell.margin_bottom = Inches(0.03)
                # 2026-05-18: anchor MIDDLE so cell text is vertically
                # centred in each row.  Without this, PowerPoint's default
                # TOP anchor placed "1,500" near the top of its row and
                # "2,000" near the top of the next row, so the ΔQ/ΔL/MPL
                # float (geometrically centred on the row BOUNDARY) sat
                # visibly below the midpoint between the two values.
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.text = str(val)
                for p in cell.text_frame.paragraphs:
                    p.alignment = PP_ALIGN.CENTER
                    for run in p.runs:
                        run.font.name = "Calibri"
                        run.font.size = Pt(16)
                        if r == 0:
                            run.font.bold = True
                            run.font.color.rgb = WHITE
                        else:
                            # Per-column color scheme (2026-05-15):
                            # L/Q black, K red, ΔL/ΔQ green, MPL blue+bold.
                            run.font.color.rgb = COL_COLORS[c]
                            if c == cols - 1:
                                run.font.bold = True
                if r == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = NAVY
                else:
                    # All data cells: white background (the cream
                    # MPL_FILL now lives on the floating MPL textboxes
                    # below, not on the cells themselves).
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = WHITE

        # ---- Floating Δ values, drawn at row boundaries ----
        # ΔL, ΔQ and MPL all float between rows; MPL gets a cream-fill
        # rounded-rect background (was the column's fill before), while
        # ΔL and ΔQ are transparent green text on the underlying white
        # cell.  Assumes equal row heights (tbl_h / rows) — true here
        # since no row has multi-line content.
        row_h = tbl_h / rows
        col_left = [tbl_left + sum(col_widths[:c]) for c in range(cols + 1)]
        float_h = Inches(0.34)                            # ~cell height
        GREEN = COL_COLORS[3]                             # 0x008000

        def _float_value(text, c, i, *, color, bold=False, fill_rgb=None,
                          border=None, line_w=0.5):
            """Place ``text`` in column c at the boundary above row i+1."""
            # 2026-05-18 (manual): per-interval Y override (hand-tuned in
            # PowerPoint).  Falls back to the mathematical row boundary
            # if no override is provided for this interval.
            boundary_y = FLOAT_CENTER_Y[i] if FLOAT_CENTER_Y[i] is not None \
                          else (tbl_top + (i + 1) * row_h)
            cell_x = col_left[c]
            cell_w = col_widths[c]
            top_y = int(boundary_y - float_h / 2)
            # Optional fill — draw a rounded rect behind the text.
            if fill_rgb is not None:
                # Inset the fill rect slightly inside the column so it
                # doesn't kiss the column-separator lines.
                pad = Inches(0.04)
                rect = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    int(cell_x + pad), top_y,
                    int(cell_w - 2 * pad), int(float_h),
                )
                try: rect.adjustments[0] = 0.18
                except Exception: pass
                rect.fill.solid()
                rect.fill.fore_color.rgb = fill_rgb
                if border is not None:
                    rect.line.color.rgb = border
                    rect.line.width = Pt(line_w)
                else:
                    rect.line.fill.background()
                rect.shadow.inherit = False
            tb = slide.shapes.add_textbox(
                int(cell_x), top_y, int(cell_w), int(float_h),
            )
            ttf = tb.text_frame
            # 2026-05-18: disable autofit — by default python-pptx writes
            # <a:spAutoFit/> on add_textbox, which causes PowerPoint to
            # shrink the textbox to its content height on render.  The
            # shrink anchors at the top, so the visible text drifts
            # upward by ~0.035" relative to the boundary line we placed
            # the box on, breaking vertical alignment with the cell text
            # in the rows above and below.
            ttf.auto_size = MSO_AUTO_SIZE.NONE
            ttf.word_wrap = True
            ttf.margin_left = Inches(0.02); ttf.margin_right = Inches(0.02)
            ttf.margin_top = Inches(0); ttf.margin_bottom = Inches(0)
            ttf.vertical_anchor = MSO_ANCHOR.MIDDLE
            pp = ttf.paragraphs[0]
            pp.alignment = PP_ALIGN.CENTER
            rr = pp.add_run()
            rr.text = text
            rr.font.name = "Calibri"
            rr.font.size = Pt(16)
            rr.font.bold = bold
            rr.font.color.rgb = color

        # Column order: c=3 → ΔQ, c=4 → ΔL, c=5 → MPL (swapped 2026-05-15).
        for i in range(1, len(L_GRID)):
            _float_value(dQ_values[i],  3, i, color=GREEN)
            _float_value(dL_values[i],  4, i, color=GREEN)
            _float_value(mpl_values[i], 5, i,
                          color=ACCENT_BLUE, bold=True,
                          fill_rgb=MPL_FILL)

        # ---- Green DOWN-arrows in the Q column, between adjacent rows ----
        # Mirrors the green connectors on the original slide 17 — visually
        # links Q[i] → Q[i+1] (the "we went from this Q to that Q" cue
        # that pairs with the ΔQ float between the same two rows).
        # Position fine-tuned on 2026-05-15: ~5 mm right of column-centre,
        # then shifted ~3 mm left so a small gap opens between the wavy
        # connector and the ΔQ digit it points at.
        q_arrow_x = (col_left[2] + int(col_widths[2] * 0.72)
                      + Inches(0.20) - Inches(0.12))
        arrow_h = Inches(0.36)                            # vertical span
        dq_col_center = col_left[3] + col_widths[3] // 2
        # Approximate width of one digit at 16 pt Calibri (used to find the
        # x-position of the FIRST digit inside a centred ΔQ value).
        char_w = Inches(0.105)
        for i in range(1, len(L_GRID)):
            # Boundary y between row i and row i+1 — pulls the hand-tuned
            # value when available so the down-arrow and wavy connector
            # line up with the floats above.
            boundary_y = FLOAT_CENTER_Y[i] if FLOAT_CENTER_Y[i] is not None \
                          else (tbl_top + (i + 1) * row_h)
            _add_arrow(slide,
                        (q_arrow_x, int(boundary_y - arrow_h / 2)),
                        (q_arrow_x, int(boundary_y + arrow_h / 2)),
                        color=GREEN, weight_pt=3.0, head=True)
            # Wavy green connector from the arrow midpoint across to the
            # ΔQ first-digit centre (shifted ~2 mm short so a small gap
            # remains between the line end and the digit).  Polyline
            # approximation of ~1.75 sine cycles for a gentle wave.
            n_chars = len(dQ_values[i])
            first_digit_x = dq_col_center - int((n_chars - 1) / 2 * char_w)
            line_end_x = first_digit_x - Inches(0.08)
            _add_wavy_line(slide,
                            q_arrow_x, line_end_x, boundary_y,
                            amplitude=Inches(0.02),
                            cycles=1.75, segments=36,
                            color=GREEN, weight_pt=1.5)

        # ---- Wide low-arc green line: first Q-arrow → Convention box ----
        # Cubic-Bezier inspired by the original slide 17.  Runs THROUGH
        # the empty horizontal band between the table header row and the
        # first row of floating ΔQ / ΔL / MPL numbers (i.e., the L = 0
        # row).  Apex sits inside that band — the curve is therefore a
        # very wide, very shallow inverted-U rather than a half-circle
        # arching over the whole table.  Stops ~0.05" before the
        # Convention box's left edge.
        arc_x_start = q_arrow_x
        arc_y_start = FLOAT_CENTER_Y[1]                     # middle of first arrow
        arc_x_end = Inches(6.10)                            # 0.05" left of conv
        arc_y_end = Inches(3.50)                            # inside the empty band
        arc_apex_y = Inches(3.35)                           # apex inside L = 0 row
        bbox_left = int(min(arc_x_start, arc_x_end))
        bbox_top = int(arc_apex_y)
        bbox_w = int(abs(arc_x_end - arc_x_start))
        bbox_h = int(max(arc_y_start, arc_y_end) - arc_apex_y)
        # Normalized coords (0–100000 along each axis of the bounding box)
        start_lx = 0
        start_ly = int(round((arc_y_start - arc_apex_y) / (max(arc_y_start, arc_y_end) - arc_apex_y) * 100000))
        end_lx = 100000
        end_ly = int(round((arc_y_end - arc_apex_y) / (max(arc_y_start, arc_y_end) - arc_apex_y) * 100000))
        # Control points pulled to the TOP and very close to the side
        # edges → inverted-U shape (steeper sides, flatter top) rather
        # than a perfectly round arc.
        cp1 = (8000, 0)
        cp2 = (92000, 0)
        weight_emu = int(1.5 * 12700)
        green_hex = f'{GREEN[0]:02X}{GREEN[1]:02X}{GREEN[2]:02X}'
        P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        A_NS_LOCAL = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        arc_xml = (
            f'<p:sp xmlns:p="{P_NS}" xmlns:a="{A_NS_LOCAL}">'
            f'<p:nvSpPr><p:cNvPr id="0" name="HalfCircleArc"/>'
            f'<p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
            f'<p:spPr>'
            f'<a:xfrm>'
            f'<a:off x="{bbox_left}" y="{bbox_top}"/>'
            f'<a:ext cx="{bbox_w}" cy="{bbox_h}"/>'
            f'</a:xfrm>'
            f'<a:custGeom>'
            f'<a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/>'
            f'<a:rect l="0" t="0" r="0" b="0"/>'
            f'<a:pathLst>'
            f'<a:path w="100000" h="100000" fill="none">'
            f'<a:moveTo><a:pt x="{start_lx}" y="{start_ly}"/></a:moveTo>'
            f'<a:cubicBezTo>'
            f'<a:pt x="{cp1[0]}" y="{cp1[1]}"/>'
            f'<a:pt x="{cp2[0]}" y="{cp2[1]}"/>'
            f'<a:pt x="{end_lx}" y="{end_ly}"/>'
            f'</a:cubicBezTo>'
            f'</a:path>'
            f'</a:pathLst>'
            f'</a:custGeom>'
            f'<a:noFill/>'
            f'<a:ln w="{weight_emu}" cap="rnd">'
            f'<a:solidFill><a:srgbClr val="{green_hex}"/></a:solidFill>'
            f'</a:ln>'
            f'</p:spPr>'
            f'</p:sp>'
        )
        slide.shapes._spTree.append(ET.fromstring(arc_xml))

        # ---- Blue arrow: bottom of MPL column → Note below the table ----
        # Visually anchors the "Note: MPL is declining" callout to the
        # MPL data above it.  Same blue as the MPL numbers (ACCENT_BLUE).
        mpl_col_center = col_left[5] + col_widths[5] // 2
        tbl_bottom = tbl_top + tbl_h
        note_top_y = Inches(6.10)        # see note positioning below
        _add_arrow(slide,
                    (mpl_col_center, int(tbl_bottom + Inches(0.10))),
                    (mpl_col_center, int(note_top_y - Inches(0.05))),
                    color=ACCENT_BLUE, weight_pt=3.0, head=True)

        # ---- Convention callout to the right of the table ----
        # 2026-05-15: narrower (5.20" → 4.20") + larger font (17 → 19 pt);
        # ΔL and ΔQ now use the same green as the body-cell ΔQ/ΔL digits.
        # Later that day: added an "Interpretation:" second paragraph
        # spelling out the first MPL value — box height bumped to 1.60"
        # to accommodate the additional 2 lines of wrapped text, then
        # widened (4.20" → 5.80") so the Interpretation line breaks
        # cleanly into exactly two lines.
        conv_w = Inches(5.80)
        conv_h = Inches(1.60)
        conv_x = tbl_left + tbl_w + Inches(0.55)
        conv_y = tbl_top + (tbl_h - conv_h) // 2     # vertically centred
        conv_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            int(conv_x), int(conv_y), int(conv_w), int(conv_h),
        )
        conv_box.fill.solid()
        conv_box.fill.fore_color.rgb = CONV_FILL
        conv_box.line.color.rgb = NAVY
        conv_box.line.width = Pt(1.0)
        conv_box.shadow.inherit = False
        try: conv_box.adjustments[0] = 0.12
        except Exception: pass
        _add_mixed_textbox(slide,
                            conv_x + Inches(0.18),
                            conv_y + Inches(0.10),
                            conv_w - Inches(0.36),
                            conv_h - Inches(0.20),
                            [
                                ('text', "Convention:  ",
                                 {'size': 19, 'bold': True, 'color': NAVY}),
                                ('text', "Compute  ",
                                 {'size': 19, 'color': NAVY}),
                                ('omml',
                                 _omml_text('Δ', color=GREEN_NUM)
                                 + _omml_run('L', color=GREEN_NUM),
                                 {'size': 19}),
                                ('text', "  and  ",
                                 {'size': 19, 'color': NAVY}),
                                ('omml',
                                 _omml_text('Δ', color=GREEN_NUM)
                                 + _omml_run('Q', color=GREEN_NUM),
                                 {'size': 19}),
                                ('text', "  for each interval",
                                 {'size': 19, 'color': NAVY}),
                                ('break', '', {}),
                                ('text', "Interpretation:  ",
                                 {'size': 19, 'bold': True, 'color': NAVY}),
                                ('text',
                                 "Between 0 and 250 workers, MPL is "
                                 "approximately 0.660",
                                 {'size': 19, 'color': NAVY}),
                            ],
                            align=PP_ALIGN.LEFT,
                            default_size=19, default_color=NAVY)

        # ---- MPL = ΔQ / ΔL formula in a cream-fill rounded-rect frame ----
        # 2026-05-18 (manual): user placed the MPL formula in PowerPoint,
        # in the empty band above the table.  Frame uses the same cream
        # FDF6E6 + navy-border styling as the Convention callout to the
        # right of the table.  Frame is drawn FIRST so the formula
        # textbox sits on top.
        # 2026-05-18 (later, manual): user nudged the formula right and
        # slightly up — from (5.876, 2.485) to (6.230, 2.393) — so it
        # sits roughly above the L=250 row of the table.
        formula_left = Inches(6.230)
        formula_top  = Inches(2.393)
        formula_w    = Inches(2.144)
        formula_h    = Inches(0.857)
        frame_pad_h  = Inches(0.10)
        frame_pad_v  = Inches(0.05)
        frame_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            int(formula_left - frame_pad_h),
            int(formula_top - frame_pad_v),
            int(formula_w + 2 * frame_pad_h),
            int(formula_h + 2 * frame_pad_v),
        )
        frame_box.fill.solid()
        frame_box.fill.fore_color.rgb = CONV_FILL
        frame_box.line.color.rgb = NAVY
        frame_box.line.width = Pt(1.0)
        frame_box.shadow.inherit = False
        try: frame_box.adjustments[0] = 0.12
        except Exception: pass
        # OMML — match slide 13's MP_L style (sSub) for cross-deck
        # consistency.
        mpl_sub = _omml_sub(_omml_run('MP'), _omml_run('L'))
        delta_q = _omml_text('Δ') + _omml_run('Q')
        delta_l = _omml_text('Δ') + _omml_run('L')
        omml_full = mpl_sub + _omml_text(' = ') + _omml_frac(delta_q, delta_l)
        _add_math_equation(slide,
                           left=formula_left, top=formula_top,
                           width=formula_w, height=formula_h,
                           omml_content=omml_full,
                           size_pt=24, color=ACCENT_BLUE)

        # ---- Blue connector: "0.660" cell → MPL = ΔQ/ΔL formula ----
        # 2026-05-18 (manual request): visually links the first MPL value
        # in the table to the MPL = ΔQ/ΔL formula above it.  Same
        # ACCENT_BLUE as the MPL column values and the formula text.
        # 2026-05-18 (later, manual): user nudged both endpoints — start
        # moved from cell-centre (5.125, 3.525) to a point 0.345" right
        # of centre; end moved from formula-centre (7.302, 2.821) to the
        # lower-left corner area of the formula box (6.330, 2.980).
        # The line is shorter and "points at" the formula rather than
        # going to its centre.
        _add_arrow(slide,
                    (Inches(5.470), Inches(3.525)),
                    (Inches(6.330), Inches(2.980)),
                    color=ACCENT_BLUE, weight_pt=1.5, head=False)

        # ---- "MPL is declining as we add workers" — Convention-style box ----
        # 2026-05-15: the Note now lives inside the same cream Convention
        # callout chrome used elsewhere on the slide.  Centred horizontally
        # at y=6.10 (the blue MPL→Note arrow above still terminates just
        # before this box).
        note_w = Inches(8.20)
        note_h = Inches(0.75)
        note_x = (SLIDE_W - note_w) // 2
        note_y = Inches(6.10)
        _add_convention_box(
            slide, note_x, note_y, note_w, note_h,
            runs=[
                ("Note:  ",
                 {'size': 22, 'bold': True, 'color': NAVY}),
                ("MPL ",
                 {'size': 22, 'bold': True, 'italic': True,
                  'color': ACCENT_BLUE}),
                ("is ",
                 {'size': 22, 'bold': True, 'color': NAVY}),
                ("declining",
                 {'size': 22, 'bold': True, 'italic': True,
                  'color': ACCENT_BLUE}),
                (" as we add workers",
                 {'size': 22, 'bold': True, 'color': NAVY}),
            ],
            size=22, align=PP_ALIGN.CENTER,
        )

    s = make_diagram_slide(
        prs, page_num=15,
        section_tag=SECTION_TAG_P1,
        title="Marginal Product of Labor (MPL):  Calculation",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "Same Cobb-Douglas function we used on slides 10 and 11, now with "
        "K = 100 robots fixed and L walked through {0, 250, 500, 1000, "
        "1500, 2000, 2500} workers.  L grid is denser at the start (two "
        "250-worker steps) so the steep first interval is visible before "
        "the 500-worker steps take over.  MPL falls from 0.660 cars/worker "
        "(going from 0 to 250 workers) down to 0.042 cars/worker (going "
        "from 2,000 to 2,500 workers) — that's the diminishing-marginal-"
        "product story in one column.  Convention: ΔL and ΔQ are step-"
        "by-step changes (current row minus previous row), NOT changes "
        "relative to L = 0.  This is the textbook MPL and it's what "
        "makes \"diminishing marginal product\" show up numerically."
    ))


def slide_13(prs):
    """Diminishing Marginal Product of Labor — mirrors original slide 18.

    Header + bullet wording restored to the original deck on 2026-05-15.
    Two dashed secant lines added to the LEFT chart to make MPL = slope
    of the production function visible; the "plot the slope" callout
    moved into the gap between the two charts so it reads as a bridge.
    """
    bullets = [
        ("We hold one input fixed (capital) and…", 0),
        ("Use more and more of a variable input (labor)", 0),
        ("Then total output will increase by less and less", 1),
        ("i.e., the Marginal Product of Labor (MPL) goes down", 1),
    ]

    K_FIX = 100
    BLUE = RGBColor(0x2E, 0x75, 0xB6)   # matches K=100 series on slide 11
    DASH_COLOR = GOLD                    # tangent lines (= MPL series color)

    def draw_pictures(slide):
        # Compute Q and MPL from the same Cobb-Douglas (Q = 3.155·√K·L^0.3)
        # we use on slides 11 and 12.  Fix K = 100 robots.  Use the full
        # PF_L_VALS (including the extra L=250 step) so this chart matches
        # the table on slide 11 exactly.
        L_vals = PF_L_VALS                     # 0, 250, 500, 1000, …, 5000
        Q_vals = [_pf_value(K_FIX, L) for L in L_vals]
        # MPL is the average slope over each interval.  Plot it at the
        # MIDPOINT of the interval (e.g., the MPL between L=250 and
        # L=500 — value 0.156 — appears on the chart at L=375).  This is
        # the convention used on slide 22 as well.
        MPL_L = [(L_vals[i-1] + L_vals[i]) // 2 for i in range(1, len(L_vals))]
        MPL_vals = [
            (Q_vals[i] - Q_vals[i-1]) / (L_vals[i] - L_vals[i-1])
            for i in range(1, len(L_vals))
        ]

        # ---- Chart frames ------------------------------------------------
        # 2026-05-15: chart-top moved down to 3.58" so the per-chart
        # caption can sit ABOVE the chart instead of below it.  Chart
        # heights unchanged → bottom = 6.43".
        cap_y = Inches(3.25)
        cap_h = Inches(0.30)
        left_chart_x  = Inches(0.40)
        left_chart_y  = Inches(3.58)
        left_chart_w  = Inches(5.90)
        left_chart_h  = Inches(2.85)
        right_chart_x = Inches(7.00)
        right_chart_y = Inches(3.58)
        right_chart_w = Inches(5.90)
        right_chart_h = Inches(2.85)

        # LEFT chart: Total output Q vs L
        left_chart_shape = _make_simple_line_chart(
            slide, left_chart_x, left_chart_y,
            left_chart_w, left_chart_h,
            categories=[f"{L:,}" for L in L_vals],
            values=Q_vals,
            line_color=BLUE,
            x_title="Workers (L)",
            y_title="Output (Q)",
            y_max=450, y_unit=50,
        )
        # RIGHT chart: MPL vs L — XY scatter so each data point sits at
        # the MIDPOINT of its interval (125, 375, 750, 1 250, …) while the
        # X-axis tick marks stay at the standard round-number positions
        # (0, 500, 1 000, …, 5 000).  Gold series colour matches the gold
        # tangent lines drawn on the left chart.
        right_chart_shape = _make_xy_line_chart(
            slide, right_chart_x, right_chart_y,
            right_chart_w, right_chart_h,
            series=[("MPL", list(zip(MPL_L, MPL_vals)),
                     DASH_COLOR, 'circle')],
            x_title="Workers (L)",
            y_title="MPL  (cars per worker)",
            x_min=0, x_max=5000, x_unit=500,
            y_min=0, y_max=0.45, y_unit=0.05,
            smooth=True,
        )

        # ---- Smooth both curves + pin the inner plot area to a known
        #      bounding box so the overlay tangent lines can be drawn at
        #      exact positions on the curve. -----------------------------
        INNER = ('0.15', '0.04', '0.80', '0.78')  # x, y, w, h fractions
        def _post_process_chart(chart_shape):
            chart_el = chart_shape.chart._element.find(qn('c:chart'))
            plot_el = chart_el.find(qn('c:plotArea')) if chart_el is not None else None
            if plot_el is not None:
                for old in plot_el.findall(qn('c:layout')):
                    plot_el.remove(old)
                layout = ET.Element(qn('c:layout'))
                ml = ET.SubElement(layout, qn('c:manualLayout'))
                lt = ET.SubElement(ml, qn('c:layoutTarget')); lt.set('val', 'inner')
                xm = ET.SubElement(ml, qn('c:xMode')); xm.set('val', 'edge')
                ym = ET.SubElement(ml, qn('c:yMode')); ym.set('val', 'edge')
                for tag, val in zip(('c:x', 'c:y', 'c:w', 'c:h'), INNER):
                    el = ET.SubElement(ml, qn(tag)); el.set('val', val)
                plot_el.insert(0, layout)
            # smooth=1 on every series so the curve is a smooth spline
            # through the data points instead of piecewise-linear segments.
            # NOTE: python-pptx also writes a CHART-LEVEL <c:smooth val="0">
            # directly under <c:lineChart>; PowerPoint honors that over
            # the series-level setting, so update / remove it too.
            for series in chart_shape.chart.series:
                ser_xml = series._element
                for sm in ser_xml.findall(qn('c:smooth')):
                    ser_xml.remove(sm)
                sm_el = ET.SubElement(ser_xml, qn('c:smooth'))
                sm_el.set('val', '1')
            # Chart-level smooth: live inside <c:plotArea>/<c:lineChart>
            line_chart = plot_el.find(qn('c:lineChart')) if plot_el is not None else None
            if line_chart is not None:
                for sm in line_chart.findall(qn('c:smooth')):
                    sm.set('val', '1')
        _post_process_chart(left_chart_shape)
        _post_process_chart(right_chart_shape)

        # ---- Dashed TANGENT lines on the LEFT chart ----------------------
        # Coordinates hand-tweaked in PowerPoint on 2026-05-15 against
        # the rendered smooth curve so each line visibly *kisses* the
        # production function at one point (the analytical Q = 5·√L
        # tangents drawn via _draw_tangent looked like secants because
        # PowerPoint's spline smoothing differs slightly from the true
        # √L curve).  Keep these exact endpoints — re-running the
        # analytical helper will re-introduce the visual mismatch.
        # Y-coords track the chart-top:
        #   chart_y = 3.40" (hand-edit baseline) → 3.58" (current, shift +0.18)
        # User hand-edit starts: 5.486, 4.512, 3.961, 3.570  (at chart_y=3.40)
        # → +0.18 shift gives: 5.666, 4.692, 4.141, 3.750
        # Tangent endpoints – hand-tweaked in PowerPoint on 2026-05-16
        # against the rendered curve so each line visibly kisses the
        # production function at one point.  Keep these exact endpoints;
        # re-deriving them analytically loses the visual match.
        t1_start = (Inches(1.275), Inches(5.783))
        t1_end   = (Inches(2.290), Inches(3.996))
        t2_start = (Inches(4.396), Inches(4.122))
        t2_end   = (Inches(5.896), Inches(3.870))
        # Steep / early tangent (≈ L = 1 000 region):
        _add_arrow(slide,
                    start_xy=t1_start, end_xy=t1_end,
                    color=DASH_COLOR, weight_pt=2.0,
                    head=False, dash='dash')
        # Flat / late tangent (≈ L = 4 000 region):
        _add_arrow(slide,
                    start_xy=t2_start, end_xy=t2_end,
                    color=DASH_COLOR, weight_pt=2.0,
                    head=False, dash='dash')

        # Caption ABOVE the LEFT chart — bold navy (not italic).
        _add_text(slide, left_chart_x, cap_y, left_chart_w, cap_h,
                   "Total output  (rising, flattening)",
                   size=13, bold=True, color=NAVY,
                   align=PP_ALIGN.CENTER, font="Calibri")
        # Caption ABOVE the RIGHT chart — same format as the left.
        # 2026-05-19 (manual): solid WHITE fill so the caption sits on
        # an opaque banner above the chart frame.
        # 2026-05-19 (manual, later): caption position nudged right +
        # down (was right_chart_x, cap_y = 7.00, 3.25) → (7.21, 3.40).
        # Same size.
        cap_right = _add_text(slide, Inches(7.21), Inches(3.40), right_chart_w, cap_h,
                   "Marginal Product of Labor  (declining)",
                   size=13, bold=True, color=NAVY,
                   align=PP_ALIGN.CENTER, font="Calibri")
        cap_right.fill.solid()
        cap_right.fill.fore_color.rgb = WHITE
        # Cream Convention callout — names the midpoint plotting
        # convention (data at L=125, 375, 750, …).  Position hand-
        # tweaked on 2026-05-16: box now sits INSIDE the upper-right of
        # the MPL chart (rather than above the chart frame), with a
        # short diagonal leader pointing at the L=750 data point on the
        # curve.  Same wording reused on slide 22.
        # 2026-05-19 (manual, later): callout narrowed from 3.50 → 2.786
        # so it tucks tighter into the upper-right of the chart frame.
        conv_w = Inches(2.786)
        conv_h = Inches(0.55)
        conv_x = Inches(8.997)
        conv_y = Inches(4.061)
        _add_convention_box(
            slide, conv_x, conv_y, conv_w, conv_h,
            prefix="Note:  ",
            body="MPL is plotted at the middle of each interval  (per our convention)",
            align=PP_ALIGN.CENTER, size=12,
            pad_v=Inches(0.04),
            line_spacing_pct=80,
        )

        # Short gold leader from below the Convention box's left edge
        # down-left to the L=750 data point on the MPL curve.  Endpoints
        # hand-tweaked in PowerPoint.
        _add_arrow(slide,
                    start_xy=(Inches(9.063), Inches(4.615)),
                    end_xy=(Inches(8.685), Inches(5.335)),
                    color=GOLD, weight_pt=1.5, head=True)

        # ---- "plot the slope" callout BETWEEN the two charts -------------
        # Combined (box + block arrow) centred horizontally on the gap
        # midpoint (x = 6.65).  Sizes bumped ~30 % on 2026-05-15 and a
        # soft drop shadow added to both shapes for visual weight.
        cb_w = Inches(1.04)              # 0.80 × 1.30
        cb_h = Inches(0.72)              # 0.55 × ~1.30
        arr_w = Inches(0.72)             # 0.55 × ~1.30
        arr_h = Inches(0.39)             # 0.30 × 1.30
        gap_mid_x = (left_chart_x + left_chart_w + right_chart_x) // 2
        # 0.32" (≈ 0.8 cm = 0.5 cm + 3 mm) left of the gap midpoint
        # per user requests on 2026-05-15.
        cb_x = gap_mid_x - (cb_w + arr_w) // 2 - Inches(0.32)
        cb_y = left_chart_y + (left_chart_h - cb_h) // 2     # vert. centred
        cb_shape = _add_rounded_filled_box(
            slide, left=cb_x, top=cb_y,
            width=cb_w, height=cb_h,
            label="plot the slope",
            fill=GOLD, text_color=NAVY,
            size=13, bold=True,
            corner_pct=0.18,
        )
        # Block right-arrow (MSO_SHAPE.RIGHT_ARROW) — much more visible
        # than a thin line connector.
        arr_shape = _add_arrow_shape(slide,
                                       left=cb_x + cb_w,
                                       top=cb_y + (cb_h - arr_h) // 2,
                                       width=arr_w, height=arr_h,
                                       direction="right", fill=GOLD)
        if arr_shape is not None:
            _add_drop_shadow(arr_shape)

        # Thin gold leader lines from the "plot the slope" callout to
        # each tangent.  Endpoints hand-tweaked on 2026-05-16 to land on
        # a visible point along each tangent (not the tangent's far end).
        # Connector to the EARLY (steep) tangent.
        _add_arrow(slide,
                    start_xy=(cb_x, cb_y + int(cb_h * 0.55)),
                    end_xy=(Inches(1.540), Inches(5.395)),
                    color=DASH_COLOR, weight_pt=1.0, head=False)
        # Connector to the LATE (flat) tangent.
        _add_arrow(slide,
                    start_xy=(cb_x + int(cb_w * 0.65), cb_y),
                    end_xy=(Inches(5.370), Inches(3.996)),
                    color=DASH_COLOR, weight_pt=1.0, head=False)

        # 2026-05-19 (manual, latest iteration):  the two right-side
        # annotation labels are tighter and moved further right (now
        # 2.26" wide, anchored at x=11.08), and their text is shortened
        # to "Very high MPL image" / "Very low MPL image" — they read
        # as clickable hyperlinks to the backup slides at the end of
        # the deck.  Arrows are now NAVY DASHED with new endpoints to
        # match the relocated labels.
        # 2026-05-19 (manual, latest iteration): each annotation is now
        # a two-run textbox — the leading "➤  " prefix is a plain navy
        # run (no underline) and the actual label is a separate run
        # that the post-build wiring pass tags with the slide-jump
        # hyperlink + underline.  Underline therefore starts exactly
        # at "Very", not on the arrow glyph.
        link_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            int(Inches(11.00)), int(Inches(1.55)),
            int(Inches(2.33)), int(Inches(0.90)),
        )
        link_box.fill.background()           # transparent
        link_box.line.color.rgb = NAVY
        link_box.line.width = Pt(1.25)
        link_box.shadow.inherit = False
        try: link_box.adjustments[0] = 0.15
        except Exception: pass

        def _annotation(top, label):
            box = slide.shapes.add_textbox(
                int(Inches(11.08)), int(top),
                int(Inches(2.26)), int(Inches(0.40)),
            )
            tf = box.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = 0
            tf.margin_top = tf.margin_bottom = 0
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            for text in ("➤  ", label):
                r = p.add_run()
                r.text = text
                r.font.name = "Calibri"
                r.font.size = Pt(14)
                r.font.color.rgb = NAVY
            return box
        _annotation(Inches(1.64), "Very high MPL image")
        _annotation(Inches(2.00), "Very low MPL image")

        # 2026-05-19 (manual, later): Arrow-1 start nudged right
        # (11.02 → 11.276) so it emerges from inside the new link box
        # closer to the right edge.  Arrow-2 start lifted up
        # (2.45 → 2.290) so it begins right at the link box's bottom.
        _add_arrow(slide,
                    start_xy=(Inches(11.276), Inches(1.870)),
                    end_xy=(Inches(8.12),    Inches(3.76)),
                    color=NAVY, weight_pt=1.5, head=True, dash='dash')
        _add_arrow(slide,
                    start_xy=(Inches(12.40), Inches(2.290)),
                    end_xy=(Inches(12.40), Inches(5.68)),
                    color=NAVY, weight_pt=1.5, head=True, dash='dash')

        # Bottom takeaway bar — nudged from 6.40 → 6.55 since the chart
        # frames are now ~0.33" taller (captions moved above).  Bar
        # bottom = 7.10, footer rule at 7.135 → clear by 0.035".
        _add_takeaway_bar(slide,
                           "Note:  MPL is the slope  (dQ / dL)  of the output curve",
                           top=Inches(6.55), fill=NAVY,
                           width=Inches(9.5), size=18)

    s = make_content_bulleted(
        prs, page_num=16,
        section_tag=SECTION_TAG_P1,
        title="Diminishing Marginal Product of Labor",
        bullets=bullets,
        # Tightened spacing on 2026-05-15 so all 4 bullets sit ABOVE
        # the charts that start at y = 3.25".  bullets_top raised by
        # ~0.5 cm (1.85 → 1.65) per user request so they clear the
        # figures' top edges.
        size=22, sub_size=20,
        line_spacing_pts=2, sub_line_spacing_pts=0,
        bullets_top=Inches(1.53),
        extras=draw_pictures,
    )
    _set_notes(s, (
        "The headline of this section. Diminishing MPL is a near-universal "
        "feature of short-run production: each additional worker has to "
        "share the same fixed capital, so the marginal contribution shrinks. "
        "This isn't a quirk of Rivian – it's nearly always true."
    ))


def slide_14(prs):
    """The Black Death and the return to labor.

    Layout matches the source: a half-page setup textbox on the left (the
    pre-1800 economy + 1348 question), wage-and-population chart on the right.
    """
    def draw(slide):
        # Top setup (full-width row 1) – context bullets.
        # 2026-05-16: bumped to size=26 main / sub=24; left/top fine-
        # tuned to (0.318, 1.668) per manual edit.
        top_bullets = [
            ("The (agriculture-based) economy before 1800:", 0),
            ("Land was the fixed factor; labor was variable", 1),
            ("Q = f (labor, land);  no capital", 1),
        ]
        _add_hierarchical_bullets(
            slide,
            left=Inches(0.318), top=Inches(1.668),
            width=RULE_W, height=Inches(1.550),
            items=top_bullets,
            size=26, sub_size=24, line_spacing_pts=8,
        )

        # Half-page setup textbox on the LEFT — rounded edges per the
        # Convention-callout style; larger fonts; the question on the
        # second line is in NAVY (dark blue) with a leading right-arrow.
        # Height trimmed from 3.0 → 2.407 on 2026-05-16 per manual edit.
        left_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            int(MARGIN), int(Inches(3.6)),
            int(Inches(5.800)), int(Inches(2.407)),
        )
        try:
            left_box.adjustments[0] = 0.08          # ~8 % corner radius
        except Exception:
            pass
        left_box.fill.solid()
        left_box.fill.fore_color.rgb = RGBColor(0xF4, 0xF1, 0xEA)  # warm parchment cream
        left_box.line.color.rgb = NAVY
        left_box.line.width = Pt(1.0)
        left_box.shadow.inherit = False
        tf = left_box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.25)
        tf.margin_right = Inches(0.25)
        tf.margin_top = Inches(0.2)
        tf.margin_bottom = Inches(0.2)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.LEFT
        r1 = p1.add_run()
        r1.text = "In 1348, the Black Death killed almost half the population (labor)."
        r1.font.name = "Calibri"
        r1.font.size = Pt(24)
        r1.font.color.rgb = NAVY
        # Blank line
        p_blank = tf.add_paragraph()
        # Highlighted question — leading right-arrow, navy bold
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = "→  What happened to the return to labor?"
        r2.font.name = "Calibri"
        r2.font.size = Pt(26)
        r2.font.bold = True
        r2.font.color.rgb = NAVY

        # Wages-and-population picture on the RIGHT — positions
        # hand-tweaked in PowerPoint on 2026-05-16: image moved up,
        # caption now sits AT THE TOP of the image (header-style)
        # instead of below.
        _add_source_image(slide, 14, "rId3",
                           left=Inches(6.653), top=Inches(2.918),
                           width=Inches(6.300))
        _add_text(slide, Inches(6.700), Inches(2.950), Inches(6.300), Inches(0.250),
                   "Wages and population, England 1300-1500",
                   size=13, italic=True, bold=True, color=NAVY,
                   align=PP_ALIGN.CENTER, font="Calibri")

        # ---- Two label-arrows on the chart (Population, Return to labor) ----
        # Positions hand-tweaked in PowerPoint on 2026-05-16.  Colors
        # differentiated per user request: Population in dark grey,
        # Return-to-labor in black — so each label/arrow visually matches
        # the chart curve it points at.
        DARK_GRAY = RGBColor(0x40, 0x40, 0x40)
        BLACK = RGBColor(0x00, 0x00, 0x00)
        # "Population" label + up-left arrow into the population curve.
        _add_text(slide, Inches(9.803), Inches(5.362),
                   Inches(0.950), Inches(0.300),
                   "Population",
                   size=12, bold=True, italic=True, color=DARK_GRAY,
                   align=PP_ALIGN.CENTER, font="Calibri")
        _add_arrow(slide,
                    start_xy=(Inches(9.850), Inches(5.425)),
                    end_xy=(Inches(9.463), Inches(5.020)),
                    color=DARK_GRAY, weight_pt=2.0, head=True)
        # "Return to labor" label + up-left arrow into the wage curve.
        _add_text(slide, Inches(8.978), Inches(5.707),
                   Inches(1.300), Inches(0.300),
                   "Return to labor",
                   size=12, bold=True, italic=True, color=BLACK,
                   align=PP_ALIGN.CENTER, font="Calibri")
        _add_arrow(slide,
                    start_xy=(Inches(9.000), Inches(5.770)),
                    end_xy=(Inches(8.620), Inches(5.223)),
                    color=BLACK, weight_pt=2.0, head=True)

    s = make_diagram_slide(
        prs, page_num=17,
        section_tag=SECTION_TAG_P1,
        title="Famous Example for Diminishing Marginal Returns",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "A favorite historical example. The Black Death killed roughly 40% "
        "of Europe's labor force in the 14th century. Wages of survivors "
        "rose sharply – consistent with their (now higher) marginal product. "
        "Real-world proof of marginal-product reasoning, 600 years before "
        "economists named it."
    ))


def slide_15(prs):
    """Rivian hiring scenario: ~$80k R1, ~$40k materials, fixed K.

    Replaces the two Tesla photos (car carrier + Model 3 EMF) with a
    single Rivian R1 photo at the right, bullets at the left.  Layout
    matches the rest of the deck: hero image on one side, narrative on
    the other, gold takeaway bar at the bottom.
    """
    def draw(slide):
        # Bullets on the LEFT  (materials cost dropped per user decision —
        # we'll handle materials separately in the cost-side of the module).
        # 2026-05-16: bumped to size=26 / sub=22 per user feedback (EMBA
        # readability — sub-bullets at 18pt were too small).
        # 2026-05-18 (manual): dropped the level-2 "(average transaction
        # price, 2024–25)" footnote; replaced with a level-1 line that
        # pivots straight from price to net revenue, paralleling the
        # framing on slide 19.
        bullets = [
            ("Demand and output price are given", 0),
            ("Large number of R1 ordered at price of ~$80k", 1),
            ("Out of this, ~$35k is material cost → (Net) Revenue per car ~$45k", 1),
            ("Short run:  capital (factory size, robots) is fixed", 0),
            ("The only way to expand production is to hire more workers", 0),
        ]
        _add_hierarchical_bullets(
            slide,
            left=MARGIN, top=Inches(1.85),
            width=Inches(8.0), height=Inches(4.4),
            items=bullets,
            size=26, sub_size=22, line_spacing_pts=8,
        )

        # Rivian R1 picture on the RIGHT (replaces the Tesla car carrier).
        rivian = OUT_DIR / "_rivian.jpg"
        if rivian.exists():
            pic = slide.shapes.add_picture(
                str(rivian),
                int(Inches(8.55)), int(Inches(1.95)),
                width=int(Inches(4.30)), height=int(Inches(3.0)),
            )
            _apply_picture_style(pic)
            # Small attribution
            _add_text(slide, Inches(8.55), Inches(5.05),
                       Inches(4.30), Inches(0.20),
                       "Rivian R1  (CC BY-SA, Wikimedia)",
                       size=9, italic=True, color=GRAY,
                       align=PP_ALIGN.CENTER, font="Calibri")

        # Bottom: rounded gold question box with drop shadow.  Narrower
        # than a full takeaway bar; leading "→ " arrow prefix anchors
        # the visual emphasis at the start of the sentence.
        # 2026-05-18 (manual): user shrank the box and slid it up + right
        # so it tucks just below the bullets rather than spanning the
        # full footer band.  Prior values: w=Inches(8.5), h=Inches(0.65),
        # centred horizontally at top=Inches(6.45).
        box_w = Inches(6.973)
        box_h = Inches(0.605)
        box_x = Inches(3.439)
        box_y = Inches(5.995)
        _add_rounded_filled_box(
            slide, box_x, box_y, box_w, box_h,
            label="→  How many workers should Rivian optimally hire?",
            fill=GOLD, text_color=NAVY,
            size=20, bold=True,
            corner_pct=0.20, shadow=True,
        )

    s = make_diagram_slide(
        prs, page_num=18,
        section_tag=SECTION_TAG_P1,
        title="Hiring Decisions in the Short Run —  Context & Scenario",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "Setup for the next several slides.  We'll derive Rivian's optimal "
        "hiring level.  Given: Rivian R1 sells at ~$80K (average "
        "transaction price reported by Rivian for 2024–25), and capital "
        "(the plant and robot count) is fixed in the short run.  Materials "
        "and other variable costs are deliberately set aside here – we "
        "handle them on the cost side of the module.  Question: how many "
        "workers should Rivian hire?"
    ))


def slide_16(prs):
    """MRPL concept – merged from source slides 17 and 18.

    Establishes that this is a SHORT-RUN concept (K fixed) and gives the
    proper textbook definition MRPL = MR × MPL.  For a price-taker firm
    MR ≈ P, so MRPL ≈ P × MPL.  Materials are NOT netted out here – they
    belong on the cost side of the module (user decision: drop materials
    from the MRPL framing to avoid confusing MC = marginal cost).
    """
    def _add_styled_box(slide, left, top, width, height, *,
                          label, fill, text_color, size, corner_adj=0.06,
                          shadow_alpha=50000):
        """Filled, slightly-rounded box with a soft drop shadow.

        corner_adj small (0.05-0.08) gives a barely-rounded rectangle;
        shadow_alpha=50000 is 50% black at 4pt blur, 3pt offset.
        """
        left, top, width, height = int(left), int(top), int(width), int(height)
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       left, top, width, height)
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
        shp.line.fill.background()
        shp.shadow.inherit = False
        try: shp.adjustments[0] = corner_adj
        except Exception: pass
        # Drop shadow via XML
        spPr = shp._element.spPr
        for old in spPr.findall(qn('a:effectLst')):
            spPr.remove(old)
        effectLst = ET.SubElement(spPr, qn('a:effectLst'))
        outerShdw = ET.SubElement(effectLst, qn('a:outerShdw'))
        outerShdw.set('blurRad', '50800')
        outerShdw.set('dist', '38100')
        outerShdw.set('dir', '2700000')
        outerShdw.set('algn', 'tl')
        outerShdw.set('rotWithShape', '0')
        rgb = ET.SubElement(outerShdw, qn('a:srgbClr'))
        rgb.set('val', '000000')
        alpha = ET.SubElement(rgb, qn('a:alpha'))
        alpha.set('val', str(int(shadow_alpha)))
        # Label centred
        tf = shp.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.15)
        tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.name = "Calibri"
        r.font.size = Pt(size)
        r.font.bold = True
        r.font.color.rgb = text_color
        return shp

    def _add_shadow(shp, alpha=50000):
        spPr = shp._element.spPr
        for old in spPr.findall(qn('a:effectLst')):
            spPr.remove(old)
        effectLst = ET.SubElement(spPr, qn('a:effectLst'))
        outerShdw = ET.SubElement(effectLst, qn('a:outerShdw'))
        outerShdw.set('blurRad', '50800')
        outerShdw.set('dist', '38100')
        outerShdw.set('dir', '2700000')
        outerShdw.set('algn', 'tl')
        outerShdw.set('rotWithShape', '0')
        rgb = ET.SubElement(outerShdw, qn('a:srgbClr'))
        rgb.set('val', '000000')
        a = ET.SubElement(rgb, qn('a:alpha'))
        a.set('val', str(int(alpha)))

    def draw(slide):
        # Short-run framing (italic navy) — sets the scope.
        # 2026-05-16: top moved up from 1.85 → 1.579 per manual edit.
        _add_text(slide, MARGIN, Inches(1.579), RULE_W, Inches(0.35),
                   "In the short run  (capital K fixed):",
                   size=18, italic=True, bold=True, color=NAVY,
                   align=PP_ALIGN.CENTER, font="Calibri")

        # ---- HERO box (two-line): name on top, plain-English below ----
        # Position + height tightened on 2026-05-16 per manual edit.
        hero_x = Inches(0.976); hero_y = Inches(1.990)
        hero_w = Inches(11.300); hero_h = Inches(1.110)
        hero = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        int(hero_x), int(hero_y),
                                        int(hero_w), int(hero_h))
        hero.fill.solid(); hero.fill.fore_color.rgb = NAVY
        hero.line.fill.background()
        hero.shadow.inherit = False
        try: hero.adjustments[0] = 0.06
        except Exception: pass
        _add_shadow(hero)
        htf = hero.text_frame
        htf.vertical_anchor = MSO_ANCHOR.MIDDLE
        htf.margin_left = Inches(0.15); htf.margin_right = Inches(0.15)
        htf.margin_top = Inches(0.05);  htf.margin_bottom = Inches(0.05)
        htf.word_wrap = True
        # Concept-accent blue for the concept name (per course CLAUDE.md);
        # remaining text white bold.
        CONCEPT_ACCENT = RGBColor(0x9E, 0xC5, 0xF7)   # soft light-blue
        p1 = htf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r0 = p1.add_run(); r0.text = "MRPL  =  "
        r0.font.name = "Calibri"; r0.font.size = Pt(28); r0.font.bold = True
        r0.font.color.rgb = WHITE
        r1 = p1.add_run(); r1.text = "Marginal Revenue Product of Labor"
        r1.font.name = "Calibri"; r1.font.size = Pt(28); r1.font.bold = True
        r1.font.color.rgb = CONCEPT_ACCENT
        p2 = htf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(4)
        r = p2.add_run()
        r.text = "the extra revenue from one more worker"
        r.font.name = "Calibri"; r.font.size = Pt(17); r.font.bold = False
        r.font.italic = True
        r.font.color.rgb = WHITE

        # ---- DECOMPOSITION box right below the HERO ----
        # Header → MR/MPL definitions → three bullets → italic note.
        # 2026-05-17: moved up + made taller per manual edit so the new
        # italic note at the bottom doesn't crowd the bullets.
        dec_x = Inches(1.000); dec_y = Inches(3.289)
        dec_w = Inches(11.200); dec_h = Inches(2.500)
        dec = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       int(dec_x), int(dec_y),
                                       int(dec_w), int(dec_h))
        dec.fill.solid(); dec.fill.fore_color.rgb = RGBColor(0xFD, 0xF6, 0xE6)
        dec.line.color.rgb = NAVY
        dec.line.width = Pt(0.75)
        dec.shadow.inherit = False
        try: dec.adjustments[0] = 0.06
        except Exception: pass
        _add_shadow(dec)
        dtf = dec.text_frame
        dtf.vertical_anchor = MSO_ANCHOR.TOP
        dtf.margin_left = Inches(0.30); dtf.margin_right = Inches(0.30)
        dtf.margin_top = Inches(0.12);  dtf.margin_bottom = Inches(0.08)
        dtf.word_wrap = True
        # Header
        ph = dtf.paragraphs[0]
        ph.alignment = PP_ALIGN.CENTER
        rh = ph.add_run()
        rh.text = "Decomposition:   MRPL  =  MR × MPL"
        rh.font.name = "Calibri"; rh.font.size = Pt(18); rh.font.bold = True
        rh.font.color.rgb = NAVY
        # MR definition (centred, italic, smaller)
        pmr = dtf.add_paragraph()
        pmr.alignment = PP_ALIGN.CENTER
        pmr.space_before = Pt(4)
        rmr = pmr.add_run()
        rmr.text = "MR:  marginal revenue from selling an extra item"
        rmr.font.name = "Calibri"; rmr.font.size = Pt(14)
        rmr.font.italic = True
        rmr.font.color.rgb = NAVY
        # MPL definition (centred) — space_before 2 → 0 on 2026-05-17 to
        # tighten the gap between MR: and MPL: lines.
        pmpl = dtf.add_paragraph()
        pmpl.alignment = PP_ALIGN.CENTER
        pmpl.space_before = Pt(0)
        rmpl = pmpl.add_run()
        rmpl.text = "MPL:  extra output (marginal product) from hiring one more worker"
        rmpl.font.name = "Calibri"; rmpl.font.size = Pt(14)
        rmpl.font.italic = True
        rmpl.font.color.rgb = NAVY
        # Bullet 1
        pb1 = dtf.add_paragraph()
        pb1.alignment = PP_ALIGN.LEFT
        pb1.space_before = Pt(10)
        rb = pb1.add_run()
        rb.text = "•  When MPL falls, MRPL falls"
        rb.font.name = "Calibri"; rb.font.size = Pt(16); rb.font.bold = True
        rb.font.color.rgb = NAVY
        # Bullet 2 — re-worded on 2026-05-17 to spell out the "less and
        # less additional output" wording.
        pb2 = dtf.add_paragraph()
        pb2.alignment = PP_ALIGN.LEFT
        pb2.space_before = Pt(8)
        rb = pb2.add_run()
        rb.text = ("•  Decreasing MPL  ⇒ marginal (additional) hires produce "
                   "less and less additional output")
        rb.font.name = "Calibri"; rb.font.size = Pt(16); rb.font.bold = True
        rb.font.color.rgb = NAVY
        # Bullet 3 — price-taker simplification (trimmed)
        pb3 = dtf.add_paragraph()
        pb3.alignment = PP_ALIGN.LEFT
        pb3.space_before = Pt(8)
        rb = pb3.add_run()
        rb.text = "•  Price-taker case:  MR = P,  so  MRPL = P × MPL"
        rb.font.name = "Calibri"; rb.font.size = Pt(16); rb.font.bold = True
        rb.font.color.rgb = NAVY
        # Italic note — added 2026-05-17.  Sits below the three bullets,
        # in italic to read as a contextual aside.
        # space_before 8 → 4 on 2026-05-17 (tighten gap to bullet 3).
        pnote = dtf.add_paragraph()
        pnote.alignment = PP_ALIGN.LEFT
        pnote.space_before = Pt(4)
        rb = pnote.add_run()
        rb.text = ("     →  Even when holding MR (or price) constant, "
                   "MRPL falls when labor is added, as a result of "
                   "falling MPL")
        rb.font.name = "Calibri"; rb.font.size = Pt(16); rb.font.italic = True
        rb.font.color.rgb = NAVY

        # ---- MB > MC anchor + DECISION RULE bar at the bottom ----
        # 2026-05-17: star bumped to (1.850 × 1.311) and shifted to
        # (10.860, 5.804) so the third line "→ Hire more" fits inside
        # the inscribed body.  Top text changed "=" → ">" — this is the
        # CONDITION for hiring more, not the optimum-state equation.
        star_w = Inches(1.850)
        star_h = Inches(1.311)
        star_x = Inches(10.860)
        star_y = Inches(5.804)
        _add_anchor_burst(
            slide, star_x, star_y, star_w, star_h,
            top_text="MB > MC",
            bottom_text="(of labor)",
            extra_text="→  Hire more",
            top_size=14, bottom_size=11,
        )

        # Gold decision-rule bar centered/left of the star.
        dr_x = Inches(3.018)
        dr_y = Inches(6.285)
        dr_w = Inches(7.216); dr_h = Inches(0.550)
        dr = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      int(dr_x), int(dr_y),
                                      int(dr_w), int(dr_h))
        dr.fill.solid(); dr.fill.fore_color.rgb = GOLD
        dr.line.fill.background()
        dr.shadow.inherit = False
        try: dr.adjustments[0] = 0.06
        except Exception: pass
        _add_shadow(dr)
        drtf = dr.text_frame
        drtf.vertical_anchor = MSO_ANCHOR.MIDDLE
        drtf.margin_left = Inches(0.15); drtf.margin_right = Inches(0.15)
        drtf.margin_top = Inches(0.05);  drtf.margin_bottom = Inches(0.05)
        drtf.word_wrap = True
        pdr = drtf.paragraphs[0]
        pdr.alignment = PP_ALIGN.CENTER
        rdr = pdr.add_run()
        rdr.text = "Decision rule:   If  MRPL > w (wage),   hire more workers"
        rdr.font.name = "Calibri"; rdr.font.size = Pt(20); rdr.font.bold = True
        rdr.font.color.rgb = NAVY

        # Arrow points FROM the star's left edge TO the rule bar's right
        # edge — leftward, ~horizontal.  Endpoints hand-tweaked.
        _add_arrow(slide,
                    start_xy=(Inches(11.026), Inches(6.465)),
                    end_xy=(Inches(10.233), Inches(6.560)),
                    color=GOLD, weight_pt=2.0, head=True)

    s = make_diagram_slide(
        prs, page_num=19,
        section_tag=SECTION_TAG_P1,
        title="Hiring Decisions in the Short Run —  Core Concept",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "MRPL — Marginal Revenue Product of Labor.  In plain terms: how "
        "much extra revenue does one more worker produce?  The textbook "
        "definition is MR × MPL — marginal revenue per unit times the "
        "marginal product of labor.  For a firm small enough that one "
        "more truck doesn't move the market price (price-taker), MR ≈ P, "
        "so MRPL ≈ P × MPL.  Since MPL falls as L grows (diminishing "
        "returns), MRPL also falls — which is why every firm has a finite "
        "optimal hiring level.  We're staying in the short run here: "
        "capital K is fixed, so the only lever is L.  Materials and other "
        "variable costs come back in the cost-side of the module."
    ))


def slide_17(prs):
    """MRPL – detail."""
    def draw(slide):
        # Major-concept formula highlighted at top – proper OMML rendering.
        # MRPL = MPL × (P − MC)  (with MRPL, MPL, MC upright acronyms; P, MC
        # italic-variable styling handled by Cambria Math).
        formula = (
            _omml_text('MRPL') + _omml_text(' = ') +
            _omml_text('MPL') + _omml_text(' × ') +
            _omml_text('(') + _omml_run('P') + _omml_text(' − ') +
            _omml_text('MC') + _omml_text(')')
        )
        _add_math_equation(
            slide,
            left=Inches(1.5), top=Inches(2.0),
            width=Inches(10.3), height=Inches(1.0),
            omml_content=formula,
            size_pt=36, color=NAVY,
        )
        _add_text(slide, MARGIN, Inches(3.05), RULE_W, Inches(0.3),
                   "(net of materials cost per unit)",
                   size=16, italic=True, color=GRAY,
                   align=PP_ALIGN.CENTER, font="Calibri")

        # Detail bullets
        bullets = [
            ("When MPL falls, MRPL falls", 0),
            ("Even if price stays constant, each additional worker is worth less", 1),
            ("The economic value of a marginal hire shrinks as you scale up", 0),
            ("Implication: there is a finite optimal number of workers", 1),
        ]
        _add_hierarchical_bullets(
            slide,
            left=MARGIN, top=Inches(3.4),
            width=RULE_W, height=Inches(3.0),
            items=bullets,
            size=26, sub_size=22, line_spacing_pts=14,
        )

        # Bottom takeaway: the optimal hiring rule preview
        _add_takeaway_bar(slide,
                           "Optimal hiring stops where MRPL just covers the wage",
                           top=Inches(6.5), fill=GOLD, text_color=NAVY,
                           width=Inches(10.0))

    s = make_diagram_slide(
        prs, page_num=20,
        section_tag=SECTION_TAG_P1,
        title="MRPL – Detail",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "Same definition, more carefully. Notice: when MPL falls, MRPL "
        "falls. So even if the price stays constant, each additional "
        "worker is worth less than the previous one. The economic value of "
        "a marginal hire shrinks as you scale up – which is why every firm "
        "has a finite optimal hiring level."
    ))


def slide_18(prs):
    """Example: MRPL at 2,300 employees and 100 robots.

    L = 2,300 is chosen deliberately so students must identify which
    table interval contains it (2,000 → 2,500) before computing MPL —
    consistent with the slide-14 MPL convention (always compute ΔQ/ΔL
    over a full table interval).
    """
    def draw(slide):
        # 2026-05-17: bullet fonts bumped to 24 / 22; the formula line
        # rewritten per manual edit — "MRPL = MR × MPL" sub-bullet
        # removed; the price-taker line tightened.
        # 2026-05-17 (later): user hand-edited slide 19 to (a) bump
        # employee count 4,000 → 4,200 so students must locate the
        # interval, and (b) reframe per-car return as Net Revenue
        # (Price $80k − material cost), not gross Price.
        # 2026-05-18: material cost tightened from ~$50k → ~$35k for
        # realism, so Net Revenue per car ~$45k (was ~$30k).
        # 2026-05-18 (later): current workforce moved 4,200 → 2,300 so
        # students locate the 2,000→2,500 interval (was 4,000→4,500).
        # 2026-05-18 (later, manual): user consolidated the price+materials
        # lines from slide 19.  The detailed "Price $80k, of which $35k
        # material cost" framing now lives on slide 17; slide 19 just
        # states the net-revenue value used in the calculation below.
        bullets = [
            ("Currently 100 robots and 2,300 employees on the R1 line", 0),
            ("(Net) Revenue per car ~$45k", 0),
            ("Assume that this is approx. constant", 1),
            ("", 0),  # spacer paragraph
            # 2026-05-17 (manual): bumped to 28 pt for emphasis.
            # 2026-05-18 (later, manual): bumped again to 32 pt.
            ("What is MRPL?  (in $ per worker, per week)", 0,
             {'size': 32}),
        ]
        # 2026-05-18 (manual): bullet fonts bumped to 28 / 28 (was 24 / 22)
        # — user wanted top-level and sub-bullets at the same 28 pt size
        # for this scenario slide.
        _add_hierarchical_bullets(
            slide,
            left=MARGIN, top=Inches(1.85),
            width=Inches(7.8), height=Inches(4.4),
            items=bullets,
            size=28, sub_size=28, line_spacing_pts=10,
        )

        # ---- Compact production-function table (same data as slide 10) ----
        # 2026-05-16: table moved RIGHT and slightly down per manual edit.
        # Axis labels disabled (with_axes=False) so we can place a larger
        # K (robots) header and a 90°-rotated L (workers) row-axis label.
        _add_compact_pf_table(slide,
                               tbl_left=Inches(9.550), tbl_top=Inches(2.013),
                               with_axes=False)
        # K axis label — wider than the default; navy italic, 16 pt
        _add_text(slide, Inches(9.374), Inches(1.689),
                   Inches(3.426), Inches(0.269),
                   "K  (robots)",
                   size=16, italic=True, color=NAVY,
                   align=PP_ALIGN.CENTER, font="Calibri")
        # L axis label — rotated 90° so it reads bottom-up like a real
        # row-axis label.
        l_box = _add_text(slide, Inches(8.720), Inches(3.728),
                           Inches(1.090), Inches(0.269),
                           "L  (workers)",
                           size=16, italic=True, color=NAVY,
                           align=PP_ALIGN.CENTER, font="Calibri")
        l_box.rotation = 270

        # Discussion-break badge — uses the helper's default top now
        # (6.25"), which sits ABOVE the gray rule at the slide bottom.
        _add_discussion_break(slide, width=Inches(4.8))

    s = make_diagram_slide(
        prs, page_num=20,
        section_tag=SECTION_TAG_P1,
        title="Example:  Calculate MRPL at 2,300 Employees and 100 Robots",
        draw_diagram=draw,
    )
    # Hyperlink the "(link)" anchor in the compact-table caption to slide 11.
    _add_slide_link_in_slide(s, "link", SLIDE_IDX_PF_TABLE, prs=prs)
    _set_notes(s, (
        "Anchor the concept with a concrete scenario. At 2,300 workers and "
        "100 robots, ask students first to find which interval in the "
        "production-function table contains L = 2,300 — that's 2,000 → 2,500 "
        "per our MPL convention.  Then reframe the per-car return: of the "
        "$80k price, about $35k is material cost, so the net revenue per "
        "car is roughly $45k.  Assume this is approximately constant across "
        "quantity, so MRPL  ≈  (Net Revenue) × MPL."
    ))


def slide_19(prs):
    """Poll: MRPL at 2,300 employees?

    Source slide is a full-bleed PollEv screenshot (single picture).  Keep
    that picture so the on-screen poll content matches what students see.
    """
    def draw(slide):
        # Render the source poll picture, centered in the body region.
        # Source pic: 9.58 x 7.08, aspect ~1.35:1.  Fit by height (5.1") so
        # width becomes ~6.9", centered horizontally.
        # 2026-05-18 (manual): picture nudged up from top=1.85 to
        # top=1.565 to leave more breathing room above the enlarged
        # bottom-right POLL pill.
        _add_source_image(slide, 19, "rId4",
                           left=Inches(3.2), top=Inches(1.565),
                           height=Inches(5.1))
        _add_text(slide, MARGIN, Inches(7.0), RULE_W, Inches(0.3),
                   "Respond at PollEv.com/nvoigtlaender",
                   size=14, italic=True, color=GRAY,
                   align=PP_ALIGN.CENTER, font="Calibri")

    s = make_diagram_slide(
        prs, page_num=21,
        section_tag=SECTION_TAG_P1,
        title="What Is Rivian's MRPL at 2,300 Employees?",
        draw_diagram=draw,
    )
    # 2026-05-18 (manual request): bottom-right poll pill with
    # discussion-break colors (gold fill + navy text).
    # 2026-05-18 (later, manual): pill enlarged to 1.487" × 0.510",
    # leading dot removed, drop shadow added — matches the visual
    # weight of other bottom-of-slide chrome elements.
    _draw_poll_pill(s, position='bottom-right',
                     fill=GOLD, text_color=NAVY, dot_color=None,
                     width=Inches(1.487), height=Inches(0.510),
                     text_size=16, shadow=True)
    _set_notes(s, (
        "Quick PollEv – compute the MRPL at 2,300 employees and submit. "
        "Give them 30 seconds. The point isn't the exact number; it's to "
        "make sure everyone identifies the right interval (2,000 → 2,500) "
        "and does the calculation in their head."
    ))


def slide_20(prs):
    """Solution: MRPL of Rivian — uses actual values from the production-
    function table.

    Applies our MPL convention: ΔQ and ΔL are computed relative to the
    previous row of the L grid.  Currently at L = 2,300 workers; the
    relevant 500-worker interval is 2,000 → 2,500.
    """
    # 2026-05-17 (manual): user restructured the solution into four
    # main-level "steps" with their derivations as sub-bullets, dropped
    # the MPL-convention hyperlink, switched the final calculation to
    # use MR (net of material cost) ≈ $30k per car → MRPL ≈ $840 per
    # worker per week (consistent with slide-19 framing).
    # 2026-05-18 (manual): emphasised the punchline line — added a
    # bold + underlined "Solution:" prefix, bolded "MRPL  =  ", and
    # changed "$840" from plain-underlined to bold + underlined.
    # (Previously: no "Solution" prefix; "MRPL  =  " plain; "$840 "
    # only underlined.)
    # 2026-05-18 (later): refreshed all numerics to match the new
    # scenario — current workforce L = 2,300 (was 4,200), interval
    # 2,000 → 2,500 (was 4,000 → 4,500), Q(2,000)=309 / Q(2,500)=330
    # (was 380 / 394), MPL = 0.042 (was 0.028), MR = $45k (was $30k),
    # final MRPL = $1,890 per worker per week (was $840).
    bullets = [
        ("Check which interval contains the current workforce  (L = 2,300)", 0),
        ("→  use the 2,000 → 2,500 step", 1,
         {'bullet_style': 'arrow', 'mar_l': 457200}),
        ("From the production-function table  (link):", 0),
        ("Q (2,000)  =  309 R1 per week", 1),
        ("Q (2,500)  =  330 R1 per week", 1),
        ("Compute MPL  =  ΔQ / ΔL", 0),
        ("MPL  =  (330 − 309) / 500  =  0.042 cars per worker per week", 1),
        ("MRPL  =  MPL × MR", 0),
        ("MR per car (net of material costs) is ≈ $45,000 ", 1),
        # Punchline: Wingdings arrow + Calibri body, with bold-underlined
        # "Solution:" prefix, bold "MRPL  =  ", and bold-underlined "$1,890".
        ([
            ('', {'wingdings': True, 'size': 20,
                         'bold': False, 'italic': False}),
            ('  ', {'size': 20, 'bold': False, 'italic': False}),
            ('Solution', {'size': 20, 'bold': True, 'italic': False,
                          'underline': True}),
            (': ', {'size': 20, 'bold': False, 'italic': False}),
            ('MRPL  =  ', {'size': 20, 'bold': True, 'italic': False}),
            ('0.042 × $45,000  =  ',
             {'size': 20, 'bold': False, 'italic': False}),
            ('$1,890', {'size': 20, 'bold': True, 'italic': False,
                        'underline': True}),
            (' per worker per week',
             {'size': 20, 'bold': False, 'italic': False}),
        ], 1, {'bullet_style': 'arrow', 'mar_l': 457200,
               'space_before_pts': 12}),
    ]
    s = make_content_bulleted(
        prs, page_num=22,
        section_tag=SECTION_TAG_P1,
        title="Solution:  MRPL of Rivian",
        bullets=bullets,
        size=24, sub_size=22, line_spacing_pts=8,
    )
    # Hyperlink the "(link)" anchor on the production-function table line.
    _add_slide_link_in_slide(s, "link", SLIDE_IDX_PF_TABLE, prs=prs)
    _set_notes(s, (
        "Reveal the answer step by step.  Step 1: ask which interval "
        "contains L = 2,300 — answer 2,000 → 2,500, because our MPL "
        "convention always computes ΔQ / ΔL over a full table interval.  "
        "Step 2: pull Q(2,000) = 309 and Q(2,500) = 330 from the "
        "production-function table at K = 100 robots.  Step 3: apply the "
        "convention — ΔQ = 21, ΔL = 500, so MPL ≈ 0.042 cars per worker "
        "per week.  Step 4: MRPL  =  MPL × MR.  MR per car is the net "
        "revenue after material cost (~$80k price − ~$35k materials = "
        "~$45k), so MRPL  ≈  0.042 × $45,000  ≈  $1,890 per worker per week. "
        "The most common slip is comparing MRPL to the TOTAL wage bill "
        "instead of the weekly wage of ONE more worker."
    ))


def slide_21(prs):
    """Hire when MRPL > wage; stop when MRPL = wage."""
    bullets = [
        ("Should Rivian hire more workers?", 0),
        ("Suppose the weekly gross wage is $1,400 per worker", 1),
        ("Yes — MRPL > wage", 1),
        ("Hiring one more worker:", 0),
        ("Revenue rises by MRPL", 1),
        ("Wage bill rises by w", 1),
        ("Profit rises whenever MRPL > w", 1),
    ]
    s = make_content_bulleted(
        prs, page_num=23,
        section_tag=SECTION_TAG_P1,
        title="Hire When MRPL > Wage;  Stop When MRPL = Wage",
        bullets=bullets,
        size=26, sub_size=22, line_spacing_pts=10,
    )
    _set_notes(s, (
        "The hiring rule in one sentence. Hire as long as the next worker "
        "brings in more than they cost. Stop the moment the next worker "
        "just breaks even. That's it – the rest is just applying this in "
        "different settings."
    ))


def slide_22(prs):
    """The optimal hiring rule in the short run (merged with old slide 21).

    Bullets on the left walk through the marginal-analysis logic;
    the chart on the right shows MRPL declining with L, the wage as
    a horizontal line, and a dashed vertical L* line at the optimal
    hiring level (L = 3,250).
    """
    # 2026-05-18 (manual): user reworked the right-hand side of the
    # slide.  Title text shortened.  Chart shrunk and shifted right
    # (was 6.5x4.3 at (6.55,1.85); now 5.23x3.36 at (7.496,1.95)) so
    # the chart's plot area expands via a manual layout (12.3%/5.8%
    # margins, 82.6%×80.1% inner area) and the legend font bumps from
    # 11pt to 14pt at a new x,y of (0.691, 0.092).  The bottom MB=MC
    # anchor + navy "Optimal Number of Workers" bar both moved up
    # (~0.5") and right (~1.3") so the bar is shorter and the burst
    # sits closer to the bar.  Bar text changed from "Optimum:" to
    # "Optimal Number of Workers:".  New dashed L* line at L=3,250
    # (added on user request).
    def draw(slide):
        # ---- Bullets (merged from slides 21 + 22) ----
        bullets = [
            ("Should Rivian hire more workers?", 0),
            ("Suppose the weekly wage  (incl. benefits)  is $1,500 per worker", 1),
            ("Hiring one more worker:", 0),
            ("Revenue rises by MRPL", 1),
            ("Wage bill rises by w", 1),
            ("Profit rises whenever MRPL > w", 1),
            ("Optimum:  hire L*  where  MRPL = w", 0),
        ]
        _add_hierarchical_bullets(
            slide,
            left=MARGIN, top=Inches(1.85),
            width=Inches(6.0), height=Inches(4.30),
            items=bullets,
            size=24, sub_size=22, line_spacing_pts=8,
        )

        # ---- Native MRPL / wage chart on the right ----
        # 2026-05-18 (morning): switched MRPL multiplier from $80 000 (gross
        # R1 price) to $30 000 (initial net-revenue framing); wage dropped
        # from $2 000 to $1 500.
        # 2026-05-18 (later): material cost tightened to ~$35k → net revenue
        # per car ~$45k.  Multiplier now $45 000; y-axis bumped to 0–$5 000
        # in steps of $500 so the top visible MRPL ($4 230 at L=750) sits
        # comfortably under the ceiling and the $1 500 wage line lands on
        # a tick.
        # MRPL = $45 000 × MPL, where MPL is the average slope over each
        # interval of PF_L_VALS at K = 100.  Plotted as an XY scatter so
        # each MRPL point sits at the MIDPOINT of its interval (same
        # convention as slide 15's MPL chart), while the X-axis tick
        # marks stay at standard L values (0, 500, 1 000, …, 5 000).
        # Skip the first two intervals (0→250 and 250→500) since their
        # MRPL ($29 700, $7 020) is off-chart; the optimal-hiring
        # intersection now lives near L ≈ 3 330 (between L=3 250/$1 530
        # and L=3 750/$1 350).
        K_FIX = 100
        L_grid = PF_L_VALS
        Q_grid = [_pf_value(K_FIX, L) for L in L_grid]
        all_mid = [(L_grid[i-1] + L_grid[i]) // 2 for i in range(1, len(L_grid))]
        all_mpl = [(Q_grid[i] - Q_grid[i-1]) / (L_grid[i] - L_grid[i-1])
                    for i in range(1, len(L_grid))]
        SKIP = 2
        mids = all_mid[SKIP:]
        mrpl_pts = [(m, int(round(45000 * mpl)))
                    for m, mpl in zip(mids, all_mpl[SKIP:])]
        WAGE = 1500
        wage_pts = [(0, WAGE), (5000, WAGE)]   # flat line across the chart

        chart_x = Inches(7.496)
        chart_y = Inches(1.950)
        chart_w = Inches(5.230)
        chart_h = Inches(3.360)
        chart_shape = _make_xy_line_chart(
            slide,
            chart_x, chart_y, chart_w, chart_h,
            series=[
                ("MRPL", mrpl_pts, NAVY, 'circle'),
                ("Wage (w)", wage_pts, GOLD, 'square'),
            ],
            x_title="L   (workers, midpoint of interval)",
            y_title="$ per worker per week",
            x_min=0, x_max=5000, x_unit=500,
            y_min=0, y_max=5000, y_unit=500,
            legend=True,
            legend_pos=('0.6888', '0.1830', '0.22', '0.20'),
            smooth=True,
        )
        # 2026-05-18 (manual): post-modify the chart to (a) bump legend
        # font to 14 pt (helper hardcodes 11 pt) and (b) add a manual
        # inner-plot-area layout so the plot fills more of the chart
        # shape.  Both values sampled from the hand-edited canonical.
        chart = chart_shape.chart
        chart.legend.font.size = Pt(14)
        plot_area = chart._chartSpace.find(qn('c:chart') + '/' + qn('c:plotArea'))
        # Remove existing <c:layout/> (auto) and insert a manualLayout
        for old in plot_area.findall(qn('c:layout')):
            plot_area.remove(old)
        layout = ET.Element(qn('c:layout'))
        ml = ET.SubElement(layout, qn('c:manualLayout'))
        ET.SubElement(ml, qn('c:layoutTarget')).set('val', 'inner')
        ET.SubElement(ml, qn('c:xMode')).set('val', 'edge')
        ET.SubElement(ml, qn('c:yMode')).set('val', 'edge')
        ET.SubElement(ml, qn('c:x')).set('val', '0.1233')
        ET.SubElement(ml, qn('c:y')).set('val', '0.0579')
        ET.SubElement(ml, qn('c:w')).set('val', '0.8264')
        ET.SubElement(ml, qn('c:h')).set('val', '0.8009')
        plot_area.insert(0, layout)

        # ---- L* dashed vertical line + label inside the chart ----
        # 2026-05-18 (manual request): add a dashed navy vertical line
        # at L = 3,250 (the optimal-hiring point — between data points
        # L=3,250/MRPL=$1,530 and L=3,750/MRPL=$1,350, where MRPL just
        # exceeds the $1,500 wage), with an "L*" label beside it.
        # 2026-05-18 (later, manual): user shortened the line so it
        # only spans from the wage line ($1,500) DOWN to the X-axis
        # (textbook "drop a vertical from the intersection" style),
        # and moved the L* label up to sit beside the lower portion
        # of the line at 16pt instead of below the X-axis at 14pt.
        L_STAR = 3250
        X_MAX  = 5000
        Y_MAX  = 5000
        WAGE_Y = 1500
        plot_x = chart_x + Inches(0.1233 * 5.230)   # left of inner plot area
        plot_y = chart_y + Inches(0.0579 * 3.360)   # top of inner plot area
        plot_w = Inches(0.8264 * 5.230)
        plot_h = Inches(0.8009 * 3.360)
        lstar_x = plot_x + int(plot_w * (L_STAR / X_MAX))
        # Line top = where the wage line crosses (= height of wage value
        # on the chart's Y-axis).  Line bottom = X-axis (plot bottom).
        wage_y_slide = plot_y + int(plot_h * (1 - WAGE_Y / Y_MAX))
        _add_arrow(slide,
                    (lstar_x, wage_y_slide),
                    (lstar_x, plot_y + plot_h),
                    color=NAVY, weight_pt=1.5, head=False, dash='dash')
        # "L*" label: navy bold italic, 16 pt, hand-positioned just to
        # the right of the line at ~75 % down its length (sampled from
        # the canonical deck).
        label_tb = slide.shapes.add_textbox(
            Inches(10.926), Inches(4.536),
            Inches(0.383), Inches(0.269),
        )
        ltf = label_tb.text_frame
        ltf.margin_left = ltf.margin_right = Inches(0)
        ltf.margin_top  = ltf.margin_bottom = Inches(0)
        ltf.word_wrap   = False
        lp = ltf.paragraphs[0]
        lp.alignment = PP_ALIGN.CENTER
        lrr = lp.add_run()
        lrr.text = "L*"
        lrr.font.name = "Calibri"
        lrr.font.size = Pt(16)
        lrr.font.bold = True
        lrr.font.italic = True
        lrr.font.color.rgb = NAVY

        # ---- Bottom: MB = MC anchor + rule statement ----
        star_w = Inches(1.6)
        star_h = Inches(1.05)
        star_x = Inches(1.583)
        star_y = Inches(5.789)
        _add_anchor_burst(
            slide, star_x, star_y, star_w, star_h,
            top_text="MB = MC",
            bottom_text="(of labor)",
            top_size=14, bottom_size=11,
        )

        # 2026-05-18 (manual request): the "Optimal Number of Workers"
        # bar gets rounded corners + a soft drop shadow.  Text has two
        # runs: the prefix at 22 pt and the rule "MRPL = w" at 24 pt
        # so the rule itself reads louder.  A colon was added after
        # "where".  Width bumped from 6.906" → 7.843" to fit the
        # larger MRPL=w portion without wrapping.
        bar_x = star_x + star_w + Inches(0.25)
        bar_y = Inches(5.975)
        bar_w = Inches(7.843)
        bar_h = Inches(0.55)
        bar_shp = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            int(bar_x), int(bar_y), int(bar_w), int(bar_h),
        )
        bar_shp.fill.solid()
        bar_shp.fill.fore_color.rgb = NAVY
        bar_shp.line.fill.background()
        try: bar_shp.adjustments[0] = 0.30
        except Exception: pass
        _add_drop_shadow(bar_shp)
        btf = bar_shp.text_frame
        btf.word_wrap = True
        btf.margin_left = btf.margin_right = Inches(0.1)
        btf.margin_top  = btf.margin_bottom = Inches(0.05)
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        bp = btf.paragraphs[0]
        bp.alignment = PP_ALIGN.CENTER
        bp_r1 = bp.add_run()
        bp_r1.text = "Optimal Number of Workers:  L*  where:  "
        bp_r1.font.name = "Calibri"
        bp_r1.font.size = Pt(22)
        bp_r1.font.bold = True
        bp_r1.font.color.rgb = WHITE
        bp_r2 = bp.add_run()
        bp_r2.text = "MRPL  =  w"
        bp_r2.font.name = "Calibri"
        bp_r2.font.size = Pt(24)
        bp_r2.font.bold = True
        bp_r2.font.color.rgb = WHITE
        _add_arrow(slide,
                    (star_x + star_w, star_y + star_h // 2),
                    (bar_x, bar_y + bar_h // 2),
                    color=GOLD, weight_pt=2.0, head=True)

    s = make_diagram_slide(
        prs, page_num=23,
        section_tag=SECTION_TAG_P1,
        title="The Optimal Hiring Rule in the Short Run",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "The same rule, in algebra: MRPL = w at the optimum. Burn this "
        "into your brain. Every short-run hiring problem you encounter as "
        "an executive is some version of this comparison."
    ))


# --------------------------------------------------------------------------
# Batch 3 – §1.1b Wage Searchers (23-29), §1.2 Long Run (30-41),
#           Part 2 section divider (42)
# --------------------------------------------------------------------------

SECTION_TAG_WAGE = "Module 3 · Production · Wage Searchers"
SECTION_TAG_LR   = "Module 3 · Production · Long Run"
SECTION_TAG_DIV  = "Module 3 · Agenda"


def slide_22b(prs):
    """Numerical solution to the optimal hiring rule (NEW slide 23).

    2026-05-18 (manual): user inserted a new slide between slide 22
    (Optimal Hiring Rule chart) and slide 23 (Wage Searchers).
    Layout is a stripped-down copy of slide 14's MPL data slide, with:
      • main table extended by a 7th MRPL column (MPL × $45 000)
      • a separate single-column "Wage" table ~1 cm to the right,
        listing six $1,500 entries vertically aligned with the
        MPL/MRPL float values on the left
      • bottom navy bar "Optimal L*  where:  MRPL  ≈  w"
    All slides after this one had their page_num bumped by 1 to
    accommodate the insertion.
    """
    K_FIX = 100
    # 2026-05-18 (later, manual): L grid extended from 2,500 → 5,000 so
    # the table shows the full range where MRPL crosses the $1,500 wage
    # (which happens between L=3,000–3,500 and L=3,500–4,000).
    L_GRID = [0, 250, 500, 1000, 1500, 2000, 2500,
              3000, 3500, 4000, 4500, 5000]
    MR = 45000          # net revenue per car ($45k, per slide 19)
    WAGE = 1500         # weekly wage per worker ($1,500, per slide 22)
    # 2026-05-18 (later, manual): user moved the table up to y=2.392
    # and let it stretch to row_h≈0.357" (same as slide 14).  Bottom
    # navy bar was deleted (the table now occupies most of the body
    # region).  Float centres land on the (now wider) row boundaries.
    # 2026-05-18 (later still): bullet, table, and wage column all
    # shifted up by ~1 cm (0.394") to free a strip of breathing room
    # below the now-larger table.
    _ROW_H = 0.357
    _TBL_TOP = 1.998
    _TBL_LEFT = 0.76
    FLOAT_CENTER_Y = [None] + [
        Inches(_TBL_TOP + (i + 1) * _ROW_H) for i in range(1, 12)
    ]
    ACCENT_BLUE = RGBColor(0x00, 0x70, 0xC0)
    BLACK_NUM  = RGBColor(0x00, 0x00, 0x00)
    RED_NUM    = RGBColor(0xC0, 0x00, 0x00)
    GREEN_NUM  = RGBColor(0x1B, 0x5E, 0x20)
    BLUE_NUM   = ACCENT_BLUE
    MPL_FILL   = RGBColor(0xFF, 0xF5, 0xE0)
    COL_COLORS = [BLACK_NUM, RED_NUM, BLACK_NUM,
                   GREEN_NUM, GREEN_NUM, BLUE_NUM, BLUE_NUM]

    def draw(slide):
        # ---- Top bullet ----
        _add_mixed_textbox(slide,
                            MARGIN, Inches(1.456),
                            RULE_W, Inches(0.45),
                            [
                                ('text', "▪  ",
                                 {'size': 24, 'bold': True, 'color': NAVY}),
                                ('text',
                                 "Compute MRPL for each L-interval",
                                 {'size': 24, 'bold': True, 'color': NAVY}),
                                ('text',
                                 "  (0-250, 250-500, 500-1,000…)",
                                 {'size': 22, 'color': NAVY}),
                            ],
                            align=PP_ALIGN.LEFT,
                            default_size=24, default_color=NAVY)

        # ---- Main 7-column table: L | K | Q | ΔQ | ΔL | MPL | MRPL ----
        col_widths = [Inches(0.80), Inches(0.65),
                       Inches(0.80), Inches(0.85),
                       Inches(0.75), Inches(0.95),
                       Inches(1.05)]
        tbl_w   = sum(col_widths)
        # 13 rows × row_h = 4.641" — fills the body band; bottom navy
        # bar removed in the user's hand-edit on 2026-05-18.
        tbl_h   = Inches(13 * _ROW_H)
        tbl_top = Inches(_TBL_TOP)
        tbl_left = Inches(_TBL_LEFT)
        cols = len(col_widths)

        Q = [_pf_value(K_FIX, L) for L in L_GRID]
        dL_values  = [None]
        dQ_values  = [None]
        mpl_values = [None]
        mrpl_values = [None]
        rows_data = [["L", "K", "Q", "ΔQ", "ΔL", "MPL", "MRPL"]]
        for i, L in enumerate(L_GRID):
            row = [f"{L:,}", f"{K_FIX}", f"{Q[i]:,}", "", "", "", ""]
            rows_data.append(row)
            if i >= 1:
                dL = L_GRID[i] - L_GRID[i-1]
                dQ = Q[i] - Q[i-1]
                mpl = dQ / dL
                mrpl = mpl * MR
                dL_values.append(f"{dL:,}")
                dQ_values.append(f"{dQ}")
                mpl_values.append(f"{mpl:.3f}")
                mrpl_values.append(f"${int(round(mrpl)):,}")
        rows = len(rows_data)

        _add_graphicframe_shadow(slide, tbl_left, tbl_top, tbl_w, tbl_h)
        tshape = slide.shapes.add_table(rows, cols, tbl_left, tbl_top,
                                          tbl_w, tbl_h)
        tbl = tshape.table
        for ci, w in enumerate(col_widths):
            tbl.columns[ci].width = w

        cell_pad_h = Inches(0.10)
        for r, row in enumerate(rows_data):
            for c, val in enumerate(row):
                cell = tbl.cell(r, c)
                cell.margin_left  = cell_pad_h
                cell.margin_right = cell_pad_h
                cell.margin_top   = Inches(0.03)
                cell.margin_bottom = Inches(0.03)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.text = str(val)
                for p in cell.text_frame.paragraphs:
                    p.alignment = PP_ALIGN.CENTER
                    for run in p.runs:
                        run.font.name = "Calibri"
                        run.font.size = Pt(16)
                        if r == 0:
                            run.font.bold = True
                            run.font.color.rgb = WHITE
                        else:
                            run.font.color.rgb = COL_COLORS[c]
                            if c >= cols - 2:   # MPL or MRPL → bold
                                run.font.bold = True
                if r == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = NAVY
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = WHITE

        # ---- Floating Δ / MPL / MRPL values ----
        col_left = [tbl_left + sum(col_widths[:c]) for c in range(cols + 1)]
        # Float height slightly less than row_h so adjacent floats have
        # a hairline gap (matches the slide_mpl_data look on slide 14).
        float_h = Inches(0.34)
        GREEN = COL_COLORS[3]

        def _float_value(text, c, i, *, color, bold=False, fill_rgb=None):
            boundary_y = FLOAT_CENTER_Y[i]
            cell_x = col_left[c]
            cell_w = col_widths[c]
            top_y = int(boundary_y - float_h / 2)
            if fill_rgb is not None:
                pad = Inches(0.04)
                rect = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    int(cell_x + pad), top_y,
                    int(cell_w - 2 * pad), int(float_h),
                )
                try: rect.adjustments[0] = 0.18
                except Exception: pass
                rect.fill.solid()
                rect.fill.fore_color.rgb = fill_rgb
                rect.line.fill.background()
                rect.shadow.inherit = False
            tb = slide.shapes.add_textbox(
                int(cell_x), top_y, int(cell_w), int(float_h),
            )
            ttf = tb.text_frame
            ttf.auto_size = MSO_AUTO_SIZE.NONE
            ttf.word_wrap = True
            ttf.margin_left = Inches(0.02); ttf.margin_right = Inches(0.02)
            ttf.margin_top  = Inches(0);    ttf.margin_bottom = Inches(0)
            ttf.vertical_anchor = MSO_ANCHOR.MIDDLE
            pp = ttf.paragraphs[0]
            pp.alignment = PP_ALIGN.CENTER
            rr = pp.add_run()
            rr.text = text
            rr.font.name = "Calibri"
            rr.font.size = Pt(16)
            rr.font.bold = bold
            rr.font.color.rgb = color

        # 2026-05-18 (later, manual): cream rounded-rect background
        # dropped from MPL floats (was distracting alongside the MRPL
        # column).  Cream fill now lives only on MRPL and on the wage
        # cells (see below).
        for i in range(1, len(L_GRID)):
            _float_value(dQ_values[i],   3, i, color=GREEN)
            _float_value(dL_values[i],   4, i, color=GREEN)
            _float_value(mpl_values[i],  5, i,
                          color=ACCENT_BLUE, bold=True)
            _float_value(mrpl_values[i], 6, i,
                          color=ACCENT_BLUE, bold=True, fill_rgb=MPL_FILL)

        # ---- Wage column to the right of the main table ----
        # 2026-05-18 (later, manual): gap widened from 1 cm to 2 cm
        # (0.788") between the main table's right edge and the wage
        # column's left edge.  Wage column now has 11 data cells (one
        # per interval) plus a backing white rect with drop shadow
        # — same chrome as the main table.
        wage_w        = Inches(1.05)
        wage_left     = tbl_left + tbl_w + Inches(0.788)
        # 2026-05-18 (later, manual): header now the same height as the
        # main table's header row (one row_h, ~0.357"); this creates a
        # visible gap between header bottom and the first wage cell
        # (which sits centred on FLOAT_CENTER_Y[1], a row+½ below
        # tbl_top).  Wage cells switched to the same cream-fill rounded
        # rect style as the MRPL floats, with ACCENT_BLUE bold text.
        header_h      = Inches(_ROW_H)
        data_cell_h   = Inches(_ROW_H)
        # Wage column backing extends from tbl_top down to the last
        # cell's bottom.
        last_cell_bot = FLOAT_CENTER_Y[11] + Inches(_ROW_H) // 2
        wage_total_h  = last_cell_bot - tbl_top

        # Backing shadow rect (same chrome as the main table)
        _add_graphicframe_shadow(slide, wage_left, tbl_top,
                                  wage_w, wage_total_h)

        # Header
        hdr = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            int(wage_left), int(tbl_top),
            int(wage_w), int(header_h),
        )
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = NAVY
        hdr.line.fill.background()
        hdr.shadow.inherit = False
        htf = hdr.text_frame
        htf.margin_left = htf.margin_right = Inches(0.05)
        htf.margin_top = htf.margin_bottom = Inches(0.03)
        htf.vertical_anchor = MSO_ANCHOR.MIDDLE
        hp = htf.paragraphs[0]
        hp.alignment = PP_ALIGN.CENTER
        hr = hp.add_run()
        hr.text = "Wage"
        hr.font.name = "Calibri"
        hr.font.size = Pt(16)
        hr.font.bold = True
        hr.font.color.rgb = WHITE

        # 11 wage cells — cream rounded rects with ACCENT_BLUE bold text.
        # Same visual treatment as the MRPL floats so the eye reads
        # "wage line" as a horizontal counterpart to MRPL.
        cell_pad_x = Inches(0.04)
        for i in range(1, 12):
            cy  = FLOAT_CENTER_Y[i]
            top = int(cy - data_cell_h // 2)
            cell = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                int(wage_left + cell_pad_x), top,
                int(wage_w - 2 * cell_pad_x), int(data_cell_h),
            )
            try: cell.adjustments[0] = 0.18
            except Exception: pass
            cell.fill.solid()
            cell.fill.fore_color.rgb = MPL_FILL
            cell.line.fill.background()
            cell.shadow.inherit = False
            ctf = cell.text_frame
            ctf.margin_left = ctf.margin_right = Inches(0.02)
            ctf.margin_top  = ctf.margin_bottom = Inches(0)
            ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
            cp = ctf.paragraphs[0]
            cp.alignment = PP_ALIGN.CENTER
            cr = cp.add_run()
            cr.text = f"${WAGE:,}"
            cr.font.name = "Calibri"
            cr.font.size = Pt(16)
            cr.font.bold = True
            cr.font.color.rgb = ACCENT_BLUE

        # 2026-05-18 (later, manual): bottom navy "Optimal L*" bar was
        # deleted by the user — the extended 13-row table now occupies
        # the full body band, so no room for it.

        # ---- Indicator arrows + ≈ sign in the gap between MRPL and
        # ---- wage columns ----
        # 2026-05-18 (manual request): the gap (~2 cm wide) between the
        # MRPL column's right edge and the wage column's left edge now
        # hosts three visual cues — a down arrow narrowing from $29,700
        # (interval 1) to $1,620 (interval 7), an up arrow narrowing
        # from $1,080 (interval 11) to $1,350 (interval 9), and a
        # centred "≈" sign at the $1,530 row (interval 8) marking
        # where MRPL approximately equals the $1,500 wage.
        gap_center_x = (col_left[7] + wage_left) // 2
        # 2026-05-18 (later, manual request): arrows thickened (2.5 →
        # 3.5 pt) and their heads pulled closer to the ≈ sign — by
        # 0.27" each, so the gap between each arrow head and the ≈
        # symbol shrinks to ~0.06" on both sides.
        # 2026-05-18 (later still, manual): arrows + ≈ sign now also
        # carry the deck's standard soft drop shadow.
        approx_cy = FLOAT_CENTER_Y[8]
        approx_half_h = Inches(0.20)        # half of approx textbox h
        arrow_gap     = Inches(0.06)        # head-to-≈ visual gap
        # Down arrow: from below $29,700 down to just above the ≈ sign.
        down_arrow = _add_arrow(slide,
                    (int(gap_center_x),
                     int(FLOAT_CENTER_Y[1] + float_h // 2)),
                    (int(gap_center_x),
                     int(approx_cy - approx_half_h - arrow_gap)),
                    color=NAVY, weight_pt=3.5, head=True)
        _add_drop_shadow(down_arrow)
        # Up arrow: from above $1,080 up to just below the ≈ sign.
        up_arrow = _add_arrow(slide,
                    (int(gap_center_x),
                     int(FLOAT_CENTER_Y[11] - float_h // 2)),
                    (int(gap_center_x),
                     int(approx_cy + approx_half_h + arrow_gap)),
                    color=NAVY, weight_pt=3.5, head=True)
        _add_drop_shadow(up_arrow)
        # "≈" sign — navy bold 28 pt — centred horizontally on the gap
        # midline and vertically on the $1,530 / $1,500 row.
        approx_w = Inches(0.40)
        approx_h = Inches(0.40)
        approx_tb = slide.shapes.add_textbox(
            int(gap_center_x - approx_w // 2),
            int(FLOAT_CENTER_Y[8] - approx_h // 2),
            int(approx_w), int(approx_h),
        )
        atf = approx_tb.text_frame
        atf.auto_size = MSO_AUTO_SIZE.NONE
        atf.word_wrap = False
        atf.margin_left = atf.margin_right = Inches(0)
        atf.margin_top  = atf.margin_bottom = Inches(0)
        atf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ap = atf.paragraphs[0]
        ap.alignment = PP_ALIGN.CENTER
        ar = ap.add_run()
        ar.text = "≈"
        ar.font.name = "Calibri"
        ar.font.size = Pt(28)
        ar.font.bold = True
        ar.font.color.rgb = NAVY
        _add_drop_shadow(approx_tb)

        # ---- Convention callout to the right of the wage column ----
        # 2026-05-18 (manual request): cream rounded-rect convention
        # box (slide-14 style) stating the optimal-hiring conclusion.
        # 2026-05-18 (later, manual): box was moved DOWN (y 3.72 →
        # 4.91, vertically aligned with the ≈/wage row) and resized
        # (4.0×1.2 → 4.492×0.673) so it sits next to the "MRPL ≈ w"
        # crossover instead of being vertically centred on the table.
        conv_left = wage_left + wage_w + Inches(0.30)
        conv_w    = Inches(4.492)
        conv_h    = Inches(0.673)
        conv_top  = Inches(4.909)
        opt_hire_box = _add_convention_box(
            slide, conv_left, conv_top, conv_w, conv_h,
            runs=[
                ("Optimal hiring",
                 {'size': 19, 'bold': True, 'color': NAVY}),
                (" falls into the interval between ",
                 {'size': 19, 'color': NAVY}),
                ("3,000", {'size': 19, 'bold': True, 'color': NAVY}),
                (" and ", {'size': 19, 'color': NAVY}),
                ("3,500", {'size': 19, 'bold': True, 'color': NAVY}),
                (" workers.", {'size': 19, 'color': NAVY}),
            ],
            size=19, align=PP_ALIGN.CENTER,
        )
        # 2026-05-18 (later, manual): drop shadow on the box so it
        # reads as a "lifted" callout next to the table.
        _add_drop_shadow(opt_hire_box)

        # ---- Connector line: wage cell at ≈ row → optimal-hiring box ----
        # 2026-05-18 (manual request): short navy line (same colour
        # and 3.5 pt weight as the arrows above) linking the $1,500
        # wage cell at the ≈ row visually to the convention callout's
        # mid-left edge.  No arrowhead — pure connector.
        # Wage cells are inset by cell_pad_x from wage_left, so the
        # cell's actual right edge = wage_left + cell_pad_x + (wage_w
        # − 2·cell_pad_x) = wage_left + wage_w − cell_pad_x.
        wage_cell_right = wage_left + wage_w - cell_pad_x
        wage_cell_cy    = FLOAT_CENTER_Y[8]
        conv_mid_y      = conv_top + conv_h // 2
        conn_line = _add_arrow(slide,
                    (int(wage_cell_right), int(wage_cell_cy)),
                    (int(conv_left),       int(conv_mid_y)),
                    color=NAVY, weight_pt=3.5, head=False)
        _add_drop_shadow(conn_line)

        # ---- "MRPL = $45,000 × MPL" formula note under the MRPL column ----
        # 2026-05-18 (manual request): small math note clarifying the
        # MRPL computation used in this slide.  Positioned just below
        # the main table, centred under the MRPL column.
        # 2026-05-18 (later, manual): user nudged the formula to the
        # right — new position (5.7765, 6.7084) sz=(2.855×0.37) — and
        # asked for a navy rounded-rect frame around it plus a
        # vertical navy line dropping from the $1,080 MRPL value
        # (last interval) down to the frame.
        mrpl_formula_omml = (
            _omml_text('MRPL') +
            _omml_text(' = $45,000 × ') +
            _omml_text('MPL')
        )
        note_left = Inches(5.7765)
        note_top  = Inches(6.7084)
        note_w    = Inches(2.855)
        note_h    = Inches(0.370)
        # Navy rounded-rect frame around the formula (slight padding so
        # the border breathes a hair).
        frame_pad = Inches(0.05)
        frame = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            int(note_left - frame_pad), int(note_top - frame_pad),
            int(note_w + 2 * frame_pad), int(note_h + 2 * frame_pad),
        )
        frame.fill.background()              # transparent fill
        frame.line.color.rgb = NAVY
        frame.line.width = Pt(1.5)
        frame.shadow.inherit = False
        try: frame.adjustments[0] = 0.15
        except Exception: pass
        # Formula on top
        _add_math_equation(
            slide,
            left=note_left, top=note_top,
            width=note_w,   height=note_h,
            omml_content=mrpl_formula_omml,
            size_pt=16, color=NAVY,
        )
        # Vertical line from below the $1,080 MRPL float (interval 11)
        # down to the top edge of the formula's navy frame.  Same navy
        # 3.5 pt style as the indicator arrows, no head.
        mrpl_col_cx = col_left[6] + col_widths[6] // 2
        line_top_y  = FLOAT_CENTER_Y[11] + float_h // 2     # below $1,080
        line_bot_y  = note_top - frame_pad                  # top of frame
        _add_arrow(slide,
                    (int(mrpl_col_cx), int(line_top_y)),
                    (int(mrpl_col_cx), int(line_bot_y)),
                    color=NAVY, weight_pt=3.5, head=False)

        # ---- MB ≈ MC anchor (12-point star) in the lower-right corner ----
        # 2026-05-18 (manual request): user introduced the deck's MB=MC
        # star pattern here too — the numerical solution IS where MB=MC
        # crystallises numerically.  Variant of the standard burst: the
        # "=" softens to "≈" because the rule lands inside an interval
        # (3,000–3,500 workers), not on a single L.  Gold star, navy
        # text, drop shadow — same chrome as every other MB=MC anchor.
        star_w = Inches(1.600)
        star_h = Inches(1.050)
        star_left = Inches(9.802)
        star_top  = Inches(5.934)
        _add_anchor_burst(
            slide, star_left, star_top, star_w, star_h,
            top_text="MB ≈ MC",
            bottom_text="(of labor)",
            top_size=14, bottom_size=11,
        )
        # Gold arrow from the star up-and-left toward the ≈ row of the
        # wage column (where MRPL ≈ w actually holds).  Head lands just
        # past the wage cell's right edge so it visually targets the
        # crossover row.
        _add_arrow(slide,
                    (int(Inches(9.990)), int(Inches(6.381))),
                    (int(Inches(8.448)), int(Inches(5.381))),
                    color=GOLD, weight_pt=2.0, head=True)

    s = make_diagram_slide(
        prs, page_num=24,
        section_tag=SECTION_TAG_P1,
        title="Optimal Hiring Rule in the Short Run:  Numerical Solution",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "Numerical solution to the optimal-hiring rule.  Re-use the "
        "MPL values from slide 14, multiply each by the net revenue "
        "per car ($45k, established on slide 19) to get MRPL.  Compare "
        "each interval's MRPL against the weekly wage of $1,500.  The "
        "optimum is the interval just BEFORE MRPL falls below the wage "
        "— here, somewhere between L=3,000 and L=3,500 if we extended "
        "the table, but for the 0–2,500 range every MRPL above $1,890 "
        "still exceeds the wage, so we'd keep hiring."
    ))


def slide_23(prs):
    """Wage searchers — merged from old slides 23 + 24.

    2026-05-18 (manual): user merged slide 24's content INTO slide 23.
    Old slide 24 ("The Case of Wage Searchers") is now this slide; the
    earlier "Caution" framing was folded into the first two bullets.
    The "wage rate is upward-sloping" line is now a sub-bullet under
    the wage-searcher Term definition.
    """
    bullets = [
        ("So far, we've assumed wages are constant", 0),
        ("Realistic for a small firm hiring at the market wage", 1),
        ("For a large firm (relative to the local labor market), hiring more workers can push the wage up", 0),
        ("Example: a local hospital hiring highly specialized surgeons", 1),
        ("Example: a frontier AI lab adding 100 senior researchers in one year", 1),
        ("Term:  the firm is a wage searcher  (not a wage taker)", 0),
        ("the wage rate is upward-sloping in employment", 1),
    ]

    def draw_extras(slide):
        # 2026-05-18 (manual request): rounded-rect + drop shadow on the
        # navy takeaway bar.  Position sampled from canonical: (2.120",
        # 6.341") sz=(9.606", 0.591").
        bar_left = Inches(2.120)
        bar_top  = Inches(6.341)
        bar_w    = Inches(9.606)
        bar_h    = Inches(0.591)
        shp = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            int(bar_left), int(bar_top), int(bar_w), int(bar_h),
        )
        shp.fill.solid()
        shp.fill.fore_color.rgb = NAVY
        shp.line.fill.background()
        try: shp.adjustments[0] = 0.30
        except Exception: pass
        _add_drop_shadow(shp)
        btf = shp.text_frame
        btf.word_wrap = True
        btf.margin_left = btf.margin_right = Inches(0.1)
        btf.margin_top = btf.margin_bottom = Inches(0.05)
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        bp = btf.paragraphs[0]
        bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run()
        br.text = "The true marginal cost of labor includes the wage-bid-up effect"
        br.font.name = "Calibri"
        br.font.size = Pt(20)
        br.font.bold = True
        br.font.color.rgb = WHITE

    s = make_content_bulleted(
        prs, page_num=25,
        section_tag=SECTION_TAG_WAGE,
        title="The Case of Wage Searchers",
        bullets=bullets,
        size=28, sub_size=24, line_spacing_pts=12,
        extras=draw_extras,
    )
    _set_notes(s, (
        "We've been assuming wages are constant — fine for a small firm "
        "hiring at the market wage.  But a large employer (relative to the "
        "local labor market) hiring many workers pushes the wage up.  Two "
        "concrete examples: a regional hospital hiring specialized surgeons "
        "can't just pay the market wage when it adds 20 of them; a frontier "
        "AI lab adding 100 senior researchers in one year drives up the "
        "going rate.  The technical term is monopsony, but you don't need "
        "the word — the intuition is enough.  Call such a firm a wage "
        "searcher (not a wage taker): the wage rate is upward-sloping in "
        "employment.  Punchline at the bottom: the TRUE marginal cost of "
        "labor includes the wage-bid-up effect — when you hire one more, "
        "you typically have to bump everyone else's wage too."
    ))


def slide_24(prs):
    """OLD slide 24 (merged into slide_23 on 2026-05-18) — kept for
    reference; not called from build_deck()."""
    bullets = [
        ("Large firm  (relative to the labor market)", 0),
        ("To recruit more labor, the firm must increase the wage", 0),
        ("For wage searchers, the wage rate is upward-sloping in employment", 0),
        ("Example: a hospital hiring highly specialized surgeons", 1),
        ("Example: Anthropic, OpenAI, DeepMind hiring senior AI researchers", 1),
    ]

    def draw_extras(slide):
        _add_takeaway_bar(slide,
                           "The true marginal cost of labor includes the wage-bid-up effect",
                           top=Inches(6.5), fill=NAVY, width=Inches(10.5))

    s = make_content_bulleted(
        prs, page_num=26,
        section_tag=SECTION_TAG_WAGE,
        title="The Case of Wage Searchers",
        bullets=bullets,
        size=26, sub_size=22, line_spacing_pts=12,
        extras=draw_extras,
    )
    _set_notes(s, (
        "The technical term is monopsony, but you don't need the word. The "
        "intuition: as a big employer hires more, the local talent pool "
        "tightens and you pay more for everyone, not just the new hire. The "
        "'true' marginal cost of labor includes this wage-bidding-up effect."
    ))


def slide_25(prs):
    """Wage-searcher caveat: equal-pay norms (chart picture)."""
    def draw(slide):
        # 2026-05-18 (manual): user shrank the embedded video frame and
        # repositioned it so the larger slide header has room to breathe.
        # Previous values: left=Inches(1.0), top=Inches(1.85),
        # width=Inches(11.3) (height auto from aspect ratio).
        _add_source_image(slide, 25, "rId5",
                           left=Inches(1.892), top=Inches(1.509),
                           width=Inches(9.546), height=Inches(5.371))
        # 2026-05-18 (manual): bottom italic caption ("Bigger firms tend
        # to pay more...") removed — the new title carries the takeaway.

    s = make_diagram_slide(
        prs, page_num=26,
        section_tag=SECTION_TAG_WAGE,
        title="Wage Searcher Caveat: Nobody likes being treated unequally…",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "Real wage data across firms of different sizes. The pattern – "
        "larger firms tend to pay more – is consistent with the wage-search "
        "story, though many other things matter too (productivity, "
        "location, benefits)."
    ))


def slide_26(prs):
    """Example: poaching an AI researcher (Anthropic + DeepMind)."""
    def draw(slide):
        bullets = [
            ("Anthropic is trying to poach a star researcher from Google DeepMind", 0),
            ("She would join Anthropic for a $5M annual salary", 1),
            ("Anthropic already employs 2 star researchers, each earning $3.5M", 0),
            ("If the new researcher is hired, the existing two will demand the same salary", 1),
        ]
        # 2026-05-18 (manual): user bumped both main and sub bullet sizes
        # to 28 pt (was 24 / 22) so the four bullets fill the body band
        # alongside the Hassabis photo.
        _add_hierarchical_bullets(
            slide,
            left=MARGIN, top=Inches(1.85),
            width=Inches(9.0), height=Inches(4.0),
            items=bullets,
            size=28, sub_size=28, line_spacing_pts=10,
        )

        # Hassabis picture (the Wikimedia / Nobel 2024 photo) on the right
        _add_source_image(slide, 26, "rId4",
                           left=Inches(9.7), top=Inches(2.0),
                           width=Inches(3.3))
        # 2026-05-18 (manual): user nudged the caption down by ~0.15" so
        # the small italic line sits clear of the picture's drop shadow.
        # Prior top: Inches(5.65).
        _add_text(slide, Inches(9.669), Inches(5.799), Inches(3.3), Inches(0.25),
                   "Demis Hassabis  (CC BY, C. Michel via Wikimedia)",
                   size=11, italic=True, color=GRAY, font="Calibri",
                   align=PP_ALIGN.CENTER)

        # Gold question box – the vote prompt for the next slide.
        # 2026-05-18 (manual): user pulled the bar up and to the left
        # (no longer centred / no longer at the footer), added a leading
        # "→" arrow, and asked for rounded corners + drop shadow so the
        # box matches the slide-17 question-bar treatment.  Prior call:
        # _add_takeaway_bar(slide, ..., top=Inches(6.45),
        # width=Inches(10.0)) — a flat full-width gold rect.
        _add_rounded_filled_box(
            slide,
            Inches(0.669), Inches(5.645),
            Inches(7.255), Inches(0.505),
            label="→  What is the marginal cost of the 3rd researcher?",
            fill=GOLD, text_color=NAVY,
            size=20, bold=True,
            corner_pct=0.20, shadow=True,
        )

        # 2026-05-18 (manual): user copied the discussion-break badge
        # into the lower-right corner — this slide kicks off a brief
        # think-through before the PollEv vote on the next slide.
        _add_discussion_break(slide, width=Inches(4.8))

    s = make_diagram_slide(
        prs, page_num=27,
        section_tag=SECTION_TAG_WAGE,
        title="Example:  The Full Cost of Poaching an AI Researcher",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "A torn-from-the-headlines wage-search example. Anthropic wants to "
        "poach a star researcher from Google DeepMind. She's rare; the third "
        "hire requires bumping up the existing senior researchers too. So "
        "the third hire is way more expensive than her salary alone. The "
        "2024-2026 AI talent wars are the textbook wage-searcher story."
    ))


def slide_27(prs):
    """Poll: what is the full marginal cost of poaching the researcher?"""
    def draw(slide):
        _add_source_image(slide, 27, "rId4",
                           left=Inches(3.2), top=Inches(1.85),
                           height=Inches(5.1))
        # 2026-05-18 (manual): "Respond at PollEv.com/nvoigtlaender"
        # caption removed — the bottom-right POLL pill (now the deck-
        # standard placement) carries the PollEv signal on its own.

    s = make_diagram_slide(
        prs, page_num=28,
        section_tag=SECTION_TAG_WAGE,
        title="What Is the Full Marginal Cost of the New Researcher?",
        draw_diagram=draw,
    )
    # 2026-05-18 (manual): switched from the small flat top-right POLL
    # pill (with leading gold dot) to the deck-standard bottom-right
    # POLL pill — gold fill, navy text, no leading dot, drop shadow,
    # enlarged to 1.487" × 0.510".  Matches the slide-20 poll chrome.
    _draw_poll_pill(s, position='bottom-right',
                     fill=GOLD, text_color=NAVY, dot_color=None,
                     width=Inches(1.487), height=Inches(0.510),
                     text_size=16, shadow=True)
    _set_notes(s, (
        "PollEv – what's the full marginal cost of the new researcher? "
        "Watch for the common trap of just reporting her $5M salary; the "
        "real answer includes the raises paid to researchers 1 and 2."
    ))


def slide_28(prs):
    """Solution: marginal cost of the 3rd researcher = $8M."""
    def draw(slide):
        # 2026-05-18 (manual): user reformatted the body — the three
        # numbered text-box "steps" and the separate navy hero box are
        # now four plain bullet points in the deck's standard "▪" style.
        # The fourth bullet is the punchline ($8M total) that the old
        # navy box used to carry on its own.
        # 2026-05-19 (manual): last bullet relabeled "Full Marginal
        # Cost" (the deck's name for this slide's quantity) and "$8M"
        # bolded to hammer the punchline.
        bullets = [
            ("The star researcher herself is paid  $5M", 0),
            ("The two existing researchers each get a raise of  ($5M − $3.5M) = $1.5M", 0),
            ("Total extra wage bill:  $5M + 2 × $1.5M", 0),
            ([
                ("Full Marginal Cost of the 3rd researcher  =  ", {}),
                ("$8M", {"bold": True}),
            ], 0),
        ]
        _add_hierarchical_bullets(
            slide,
            left=MARGIN, top=Inches(1.85),
            width=RULE_W, height=Inches(3.8),
            items=bullets,
            size=24, sub_size=22, line_spacing_pts=14,
        )

        # 2026-05-18 (manual): gold "Take-Away" bar moved up from
        # top=Inches(6.45) and switched from a flat full-width rect to
        # a rounded box with drop shadow — matches the slide-17 / 26
        # question-bar treatment.  Leading bold "Take-Away:" prefix
        # signals this is the punchline.
        # 2026-05-19 (manual): nudged up again (top 5.85 → 5.28) and a
        # touch right (left 1.811 → 1.881) so it sits cleaner under the
        # bullet block.  Wording tightened: "move the market price" →
        # "can affect wages".
        _add_rounded_filled_box(
            slide,
            Inches(1.881), Inches(5.28),
            Inches(9.711), Inches(0.55),
            label="Take-Away:  Big buyers of scarce talent can affect wages  —  factor it in",
            fill=GOLD, text_color=NAVY,
            size=20, bold=True,
            corner_pct=0.20, shadow=True,
        )

    s = make_diagram_slide(
        prs, page_num=29,
        section_tag=SECTION_TAG_WAGE,
        title="Solution:  Marginal Cost of the 3rd Researcher = $8M",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "Reveal: marginal cost of the third researcher is $8M, not her $5M "
        "salary. The lesson: when you're a big enough buyer of scarce talent "
        "(or anything), your hiring moves the market price. Factor that in. "
        "The same logic applied at Meta when they paid up to keep AI "
        "researchers from leaving in 2024."
    ))


def slide_29(prs):
    """Are real-world wages = MRPL? (UC wage-search tool)"""
    def draw(slide):
        # 2026-05-19 (manual): user reinstated the source deck's two
        # stacked screenshots — rId4 (the raw UCLA salary list) behind,
        # rId5 (the search-tool UI) on top.  Keep both: each conveys a
        # different angle on "look at the real numbers".  The back image
        # carries no shadow so the visible (top) image is the one that
        # pops off the slide.
        _add_source_image(slide, 29, "rId4",
                           left=Inches(2.516), top=Inches(1.553),
                           width=Inches(8.3), shadow=False)
        _add_source_image(slide, 29, "rId5",
                           left=Inches(2.516), top=Inches(1.545),
                           width=Inches(8.3))

        # 2026-05-19 (manual): UC search-tool URL toned down from the
        # deck-standard navy/white bold pill — now a plain transparent
        # textbox with black 14 pt regular text.  The URL is useful as
        # a pointer for students but doesn't need to dominate.
        # Position preserved from hand-edit (left 2.5 → 5.034, top 6.4
        # → 6.549) — the box hugs the right edge below the picture.
        url_box = slide.shapes.add_textbox(
            Inches(5.034), Inches(6.549),
            Inches(8.3), Inches(0.55),
        )
        url_tf = url_box.text_frame
        url_tf.word_wrap = True
        url_tf.margin_left = Inches(0.1)
        url_tf.margin_right = Inches(0.1)
        url_tf.margin_top = Inches(0.05)
        url_tf.margin_bottom = Inches(0.05)
        url_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = url_tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        # "Search tool:  " label (plain black, no hyperlink).
        r_label = p.add_run()
        r_label.text = "Search tool:  "
        r_label.font.name = "Calibri"
        r_label.font.size = Pt(14)
        r_label.font.bold = False
        r_label.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
        # 2026-05-19 (manual): URL itself is now a clickable hyperlink.
        # Keep the toned-down look (black, no underline) so the link
        # doesn't shout — click-through still works in slideshow mode.
        r_url = p.add_run()
        r_url.text = "https://ucannualwage.ucop.edu/wage/"
        r_url.font.name = "Calibri"
        r_url.font.size = Pt(14)
        r_url.font.bold = False
        r_url.font.underline = False
        r_url.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
        r_url.hyperlink.address = "https://ucannualwage.ucop.edu/wage/"

    s = make_diagram_slide(
        prs, page_num=30,
        section_tag=SECTION_TAG_WAGE,
        title="Are Real-World Wages = MRPL?",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "The empirical question that closes this section. Do real-world "
        "wages roughly equal MRPL? The UC wage-comparison tool lets you "
        "check for yourself. Spoiler: yes, broadly, but with persistent "
        "gaps that economists still argue about."
    ))


def slide_30(prs):
    """Section divider – Part 1.c: Long Run.

    Same Layout 2 / agenda view as slide 7 (Part 1 navy, Part 2 faded), but
    with an action title signalling the sub-section transition.
    """
    # 2026-05-19: renumbered from "Part 1.2" to "Part 1.c" — Part 1 now
    # has three subsections (a Production Function, b Short Run, c Long
    # Run) so this divider should match the alpha-labeled agenda.
    # 2026-05-19 (later): fade subs a and b within Part 1; only sub c
    # (Long Run) stays navy.  Matches the new Short Run agenda divider
    # on page 13.
    s = make_section_agenda(
        prs, page_num=31,
        current_part_idx=0,
        current_sub_idx=2,
        section_tag=SECTION_TAG_DIV,
        title="Part 1.c:  Long Run – Choosing the Right Input Mix",
    )
    _set_notes(s, (
        "Switching gears now from short run to long run. In the long run, "
        "capacity is no longer fixed – we get to choose K AND L from "
        "scratch. New decision: what's the right MIX of capital and labor?"
    ))


def slide_31(prs):
    """Long-run context: Rivian builds a new Georgia plant."""
    def draw(slide):
        bullets = [
            ("Context: Rivian builds its new Georgia plant", 0),
            ("Both capital and labor are flexible inputs", 0),
            ("What is the optimal input mix?", 0),
            ("E.g., robots (K) and workers (L)", 1),
            ("We will:", 0),
            ("Use marginal analysis", 1),
            ("Learn a simple rule for the optimal combination of inputs", 1),
        ]
        # 2026-05-19 (manual): narrowed bullets from Inches(8.0) → Inches(6.8)
        # so the new (wider, further-left) Rivian Georgia plant photo
        # has room on the right.
        _add_hierarchical_bullets(
            slide,
            left=MARGIN, top=Inches(1.85),
            width=Inches(6.8), height=Inches(4.6),
            items=bullets,
            size=24, sub_size=22, line_spacing_pts=10,
        )

        # 2026-05-19 (manual): user swapped the old Wikimedia "Rivian R1"
        # source-deck image for this Brian Cassella / Tribune photo of
        # the Normal, IL plant assembly line — image23.jpeg in the
        # canonical deck, persisted to _rivian_georgia.jpg in this folder.
        # Rectangular crop kept; soft drop shadow lifts it off the slide.
        rivian = OUT_DIR / "_rivian_georgia.jpg"
        if rivian.exists():
            pic = slide.shapes.add_picture(
                str(rivian),
                int(Inches(7.234)), int(Inches(1.619)),
                width=int(Inches(5.825)), height=int(Inches(3.276)),
            )
            _add_drop_shadow(pic)

        # Attribution caption — italic gray, sits just below the photo.
        _add_text(slide, Inches(8.171), Inches(4.989),
                   Inches(4.398), Inches(0.20),
                   "Brian Cassella  |  Tribune News Service  |  Getty Images",
                   size=10.5, italic=True, color=GRAY, font="Calibri",
                   align=PP_ALIGN.CENTER)

        # 2026-05-19 (manual): the old flat full-width takeaway bar is
        # now a rounded navy box sized to match the photo's width and
        # aligned under it.  Text tightened to "optimal K-and-L mix"
        # (was "the right K-and-L mix from scratch").  Rounded corners +
        # soft drop shadow — same treatment as the slide-28 take-away.
        _add_rounded_filled_box(
            slide,
            Inches(3.754), Inches(6.123),
            Inches(5.825), Inches(0.497),
            label="Long run  ⇒  pick the optimal K-and-L mix",
            fill=NAVY, text_color=WHITE,
            size=20, bold=True,
            corner_pct=0.20, shadow=True,
        )

    s = make_diagram_slide(
        prs, page_num=32,
        section_tag=SECTION_TAG_LR,
        title="Long Run:  Rivian Builds a New Georgia Plant",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "Concrete setup: Rivian building its new plant in Stanton Springs, "
        "Georgia (recently revived after the VW partnership in 2024). They "
        "get to pick everything – plant size, machinery, workforce, "
        "layout. How should they choose?"
    ))


def slide_32(prs):
    """Optimal combination of inputs (concept introduction).

    2026-05-19 (manual, late): user consolidated the old "Bang for the
    Buck Rule" chart slide (deleted slide_33) into this page — the rule
    name now lives in a rounded NAVY anchor below the formula, the
    three "when it holds" sub-bullets moved here as a bottom-left
    cluster, and a right-side legend defines the symbols.
    """
    def draw(slide):
        # 2026-05-19 (manual): top bullet pulled UP from y=1.85 to y=1.44
        # to start the body content closer to the title rule.  Size 32 pt.
        _add_hierarchical_bullets(
            slide,
            left=MARGIN, top=Inches(1.435),
            width=RULE_W, height=Inches(0.55),
            items=[("Optimal Input Mix:  How much K and how much L?", 0)],
            size=32, line_spacing_pts=0,
        )

        # 2026-05-19 (manual): "Decision rule for the long run:" — same
        # 28-pt size, but pulled UP to y=2.265 (was 2.85) along with
        # the bullet above and the formula below.
        _add_text(slide, MARGIN, Inches(2.265), RULE_W, Inches(0.5),
                   "Decision rule for the long run:",
                   size=28, italic=False, color=GRAY,
                   align=PP_ALIGN.CENTER, font="Calibri")

        # 2026-05-19 (manual): formula pulled UP to y=2.955 (was 3.54).
        _add_math_equation(
            slide,
            left=Inches(2.5), top=Inches(2.955),
            width=Inches(8.3), height=Inches(1.5),
            omml_content=_formula_bang_for_buck(),
            size_pt=44, color=NAVY,
        )

        # 2026-05-19 (manual, late): navy "Bang for the Buck Rule"
        # anchor RESHAPED — no longer centered-and-full-width.  Now a
        # smaller pill sitting in the middle of the slide between the
        # formula and the rescued sub-bullets, leaving room on the right
        # for the symbol legend.  Position (4.369, 4.759), size
        # 4.788 × 0.618".  Same NAVY/white + adj=30000 + shadow style.
        bar_x = Inches(4.369)
        bar_y = Inches(4.759)
        bar_w = Inches(4.788)
        bar_h = Inches(0.618)
        bar_shp = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            int(bar_x), int(bar_y), int(bar_w), int(bar_h),
        )
        bar_shp.fill.solid()
        bar_shp.fill.fore_color.rgb = NAVY
        bar_shp.line.fill.background()
        try: bar_shp.adjustments[0] = 0.30
        except Exception: pass
        _add_drop_shadow(bar_shp)
        btf = bar_shp.text_frame
        btf.word_wrap = True
        btf.margin_left = btf.margin_right = Inches(0.1)
        btf.margin_top  = btf.margin_bottom = Inches(0.05)
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        bp = btf.paragraphs[0]
        bp.alignment = PP_ALIGN.CENTER
        bp_r = bp.add_run()
        bp_r.text = "Bang for the Buck Rule"
        bp_r.font.name = "Calibri"
        bp_r.font.size = Pt(28)
        bp_r.font.bold = True
        bp_r.font.color.rgb = WHITE

        # 2026-05-19 (manual, late): legend container REPOSITIONED +
        # RECOLORED — was cream-fill convention-style at (10.75, 3.45);
        # now a white rounded box at (9.44, 3.12) sized 3.57 × 1.29.
        # The "Where:" header was removed (the four entries speak for
        # themselves), and MP_K / MP_L now use proper baseline-shifted
        # subscripts (italic-bold K/L) rather than Unicode ₖ ₗ.
        leg_x = Inches(9.440)
        leg_y = Inches(3.124)
        leg_w = Inches(3.568)
        leg_h = Inches(1.290)
        leg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            int(leg_x), int(leg_y), int(leg_w), int(leg_h),
        )
        leg.fill.solid()
        leg.fill.fore_color.rgb = WHITE
        leg.line.color.rgb = NAVY
        leg.line.width = Pt(1.0)
        leg.shadow.inherit = False
        try: leg.adjustments[0] = 0.10
        except Exception: pass

        # Inner textbox with 4 legend lines.
        pad_h = Inches(0.08)
        pad_v = Inches(0.025)
        ltb = slide.shapes.add_textbox(
            int(leg_x + pad_h), int(leg_y + pad_v),
            int(leg_w - 2 * pad_h), int(leg_h - 2 * pad_v),
        )
        ltf = ltb.text_frame
        ltf.word_wrap = True
        ltf.margin_left = Inches(0.05); ltf.margin_right = Inches(0.05)
        ltf.margin_top = 0; ltf.margin_bottom = 0
        ltf.vertical_anchor = MSO_ANCHOR.MIDDLE

        def _legend_line(idx, runs_spec):
            """Add a paragraph with mixed-style runs.  Each run is
            (text, opts) where opts may include bold, italic, size,
            and 'subscript' (sets baseline=-25000 on the rPr)."""
            p = ltf.paragraphs[0] if idx == 0 else ltf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            for text, opts in runs_spec:
                r = p.add_run()
                r.text = text
                r.font.name = 'Calibri'
                if 'size' in opts:
                    r.font.size = Pt(opts['size'])
                r.font.bold = opts.get('bold', False)
                r.font.italic = opts.get('italic', False)
                r.font.color.rgb = NAVY
                if opts.get('subscript'):
                    r._r.get_or_add_rPr().set('baseline', '-25000')

        VAR = 18   # default body-text size; visible in legend
        DEF = 16   # definition text size, slightly smaller
        _legend_line(0, [
            ('MP', {'bold': True, 'size': VAR}),
            ('K',  {'bold': True, 'size': VAR, 'subscript': True}),
            ('  :  Marginal Product of Capital', {'size': DEF}),
        ])
        _legend_line(1, [
            ('MP', {'bold': True, 'size': VAR}),
            ('L',  {'bold': True, 'size': VAR, 'subscript': True}),
            ('  :  Marginal Product of Labor', {'size': DEF}),
        ])
        _legend_line(2, [
            ('p',  {'bold': True, 'italic': True, 'size': VAR}),
            ('ₖ',  {'bold': True, 'size': VAR}),
            ('  :  Price of Capital', {'size': DEF}),
        ])
        _legend_line(3, [
            ('w',  {'bold': True, 'italic': True, 'size': VAR}),
            ('  :  Wage', {'size': DEF}),
        ])

        # 2026-05-19 (manual, late): three sub-bullets rescued from the
        # deleted slide_33 — they explain the conditions under which
        # the bang-for-the-buck rule holds.  Bottom-left, sub-level
        # (gray "–" markers) so they read as fine print under the
        # navy rule anchor above them.
        _add_hierarchical_bullets(
            slide,
            left=Inches(0.233), top=Inches(5.795),
            width=Inches(6.433), height=Inches(1.195),
            items=[
                ("Holds when both L and K are flexible", 1),
                ("Refers to a given output quantity Q", 1),
                ("Assumes input prices w and pₖ are constant", 1),
            ],
            size=24, sub_size=22, line_spacing_pts=8,
        )

        # 2026-05-19 (manual): teaching-note card, peach fill.  Slight
        # right-shift of x position (8.207 → 8.374) per the hand-edit.
        _add_teaching_note(
            slide,
            "Bang-for-the-Buck Rule",
            left=Inches(8.374), top=Inches(6.415),
            width=Inches(4.889), height=Inches(0.6),
            rounded=True, pdf_icon=True, label_color=NAVY,
            fill_rgb=RGBColor(0xFA, 0xC0, 0x90),
        )

    s = make_diagram_slide(
        prs, page_num=33,
        section_tag=SECTION_TAG_LR,
        title="Optimal Combination of Inputs",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "The general framework: we need a decision rule for combining "
        "inputs when both are variable. The rule will look familiar – it's "
        "the same logic as the short-run hiring rule, generalized. Three "
        "conditions under the rule: both inputs flexible, output Q held "
        "fixed, input prices w and p_K constant. Spend each additional "
        "dollar on whichever input gives the most extra output per dollar."
    ))


# 2026-05-19: slide_33 (the standalone "Bang for the Buck Rule" page)
# was removed by the user.  The rule itself now lives on slide_32's
# page; the three sub-bullets that explained when it holds were
# moved into slide_32 as well.  Function deleted to avoid a stale
# page_num=34 reference.


def slide_34(prs):
    """Applying the 'Bang for the Buck' rule — recipe for exams.

    2026-05-19: rewritten to mirror the source deck's slide 40 — two
    if-clauses naming which direction to shift dollars when the
    MP-per-$ ratios are unequal, a nested "decreasing marginal returns"
    explanation of why the adjustment converges, and a closing
    "continue until equal" line.  Replaces the prior numbered five-step
    layout.  Teaching-note card removed per user request — the
    procedure is now self-contained on the slide.
    """
    def draw(slide):
        # Subtitle directly under the title rule — italic gray small,
        # signals "this is the procedural / how-to variant".
        _add_text(slide, MARGIN, Inches(1.30), RULE_W, Inches(0.45),
                   "('recipe' for exams)",
                   size=20, italic=True, color=GRAY,
                   align=PP_ALIGN.LEFT, font="Calibri")

        # OMML formula chunks reused across bullets.  Built once so the
        # subscripted variables render in Cambria Math (matching the
        # headline formula on slide 33).
        f_mpl_w   = _formula_mp_ratio('L', 'w')   # MP_L / w   (stacked)
        f_mpk_pk  = _formula_mp_ratio('K', 'p')   # MP_K / p_K (stacked)
        f_mpl     = _omml_sub(_omml_run('MP'), _omml_run('L'))   # MP_L
        f_mpk     = _omml_sub(_omml_run('MP'), _omml_run('K'))   # MP_K
        f_eq      = f_mpl_w + _omml_text('  =  ') + f_mpk_pk     # MP_L/w = MP_K/p_K

        # Body — hierarchical bullets.  Formula portions are real inline
        # OMML math zones (run_opts={'omml': True}); prose runs stay
        # regular.  Single-letter L / K in prose stay italic.
        # 2026-05-19 (manual): user hand-tweaked spacing throughout —
        #  • first line presented as a heading (no bullet, no indent);
        #    wording "the same" → "a given";
        #  • blank spacer paragraphs removed in favour of explicit
        #    space_before_pts=18 on the three major "If" / "continue"
        #    lines (cleaner visual rhythm than empty paragraphs);
        #  • padding around `>`, `<`, and `→` reduced (was 4 spaces, now
        #    1–3) so the formulas sit tighter against the prose;
        #  • each sub-sub-bullet now closes with "as L increases" to
        #    make the diminishing-returns mechanism explicit.
        bullets = [
            ("How to adjust L and K in order to produce a given quantity at lower cost",
             0, {'bullet_style': 'arrow', 'mar_l': 0, 'indent': 0}),
            ([
                ("If  ",                 {}),
                (f_mpl_w,                {'omml': True}),
                (" > ",                  {}),
                (f_mpk_pk,               {'omml': True}),
                ("   →    use more ",    {}),
                ("L",                    {'italic': True}),
                (" and less ",           {}),
                ("K",                    {'italic': True}),
            ], 0, {'space_before_pts': 18}),
            ("Decreasing marginal returns imply:", 1),
            ([
                ("Due to increased ",    {}),
                ("L",                    {'italic': True}),
                (",  ",                  {}),
                (f_mpl,                  {'omml': True}),
                ("   will fall as ",     {}),
                ("L",                    {'italic': True}),
                (" increases",           {}),
            ], 2),
            ([
                ("Due to reduced ",      {}),
                ("K",                    {'italic': True}),
                (",  ",                  {}),
                (f_mpk,                  {'omml': True}),
                ("   will rise as ",     {}),
                ("L",                    {'italic': True}),
                (" increases",           {}),
            ], 2),
            ([
                ("If  ",                 {}),
                (f_mpl_w,                {'omml': True}),
                (" < ",                  {}),
                (f_mpk_pk,               {'omml': True}),
                ("   →    use less ",    {}),
                ("L",                    {'italic': True}),
                (" and more ",           {}),
                ("K",                    {'italic': True}),
            ], 0, {'space_before_pts': 18}),
            ([
                ("…continue to adjust until    ", {}),
                (f_eq,                            {'omml': True}),
            ], 0, {'space_before_pts': 18}),
        ]
        _add_hierarchical_bullets(
            slide,
            left=MARGIN, top=Inches(1.95),
            width=RULE_W, height=Inches(5.0),
            items=bullets,
            size=28, sub_size=22, line_spacing_pts=10,
        )

    s = make_diagram_slide(
        prs, page_num=34,
        section_tag=SECTION_TAG_LR,
        title="Applying the 'Bang for the Buck' Rule",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "Recipe for applying the rule.  Two cases when the MP-per-dollar "
        "ratios are not yet equal:  if MP_L/w is the higher one, the firm "
        "should shift toward labor (use more L, less K);  if MP_K/p_K is "
        "higher, shift toward capital.  As the adjustment happens, "
        "diminishing marginal returns kick in — adding L pulls MP_L down, "
        "taking K away pushes MP_K up — so the two ratios converge.  "
        "Keep going until they meet:  MP_L/w = MP_K/p_K is the optimum."
    ))


def slide_35(prs):
    """Example: Rivian's New Georgia plant.

    2026-05-19 (manual round 3): operating point reverted from the
    off-grid K = 320 / L = 1,800 back to grid-aligned K = 300 / L = 2,000
    — the nearest-grid-points convention added in round 2 was deemed too
    complex for an intro to the bang-for-the-buck rule, so slides 36–39
    return to direct table lookup.  Q line stays dropped (user preference
    from round 2 hand-edits).  Bottom takeaway content remains folded
    into the last bullet on the slide.
    """
    def draw(slide):
        bullets = [
            ("Rivian is building a new plant in Stanton Springs, Georgia", 0),
            ("They ask for your advice on the optimal mix of robots and workers", 0),
            ("You know:", 0),
            ("Current plan:  300 robots and 2,000 workers", 1),
            ("Weekly wage for suitable workers:  w = $1,200", 1),
            ("Cost of one robot (per week):  pₖ = $20,000", 1),
            ("Is Rivian's 300 robots / 2,000 workers plan optimal?", 0),
        ]
        _add_hierarchical_bullets(
            slide,
            left=MARGIN, top=Inches(1.85),
            width=Inches(8.0), height=Inches(4.4),
            items=bullets,
            size=24, sub_size=22, line_spacing_pts=8,
        )

        # 2026-05-19 (manual): replaced the Wikimedia R1 photo with a
        # new project rendering for the Georgia plant.  Apply rounded
        # corners + soft drop shadow so the picture sits cleanly.
        # Position + size hand-set by user; caption dropped.
        rivian = OUT_DIR / "_rivian_georgia_plan.png"
        if rivian.exists():
            pic = slide.shapes.add_picture(
                str(rivian),
                int(Inches(8.051)), int(Inches(1.405)),
                width=int(Inches(5.187)), height=int(Inches(3.460)),
            )
            _apply_picture_style(pic)

        # 2026-05-19 (manual round 5): user added a Discussion Break badge
        # at the bottom-right.  Captured from the canonical deck.
        _add_discussion_break(slide)

    s = make_diagram_slide(
        prs, page_num=35,
        section_tag=SECTION_TAG_LR,
        title="Example:  Rivian's New Georgia Plant",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "Rivian announced a second US assembly plant in Stanton Springs, "
        "Georgia – a multi-billion-dollar project meant to add capacity "
        "once it ramps.  We'll apply the bang-for-the-buck rule to a "
        "stylised version of those plans:  robots vs. workers, given each "
        "input's price and marginal product.  Current plan is K = 300 "
        "robots and L = 2,000 workers — both on the production-function "
        "table grid, so the next slides read Q, MPL, MPK straight off the "
        "table.  The strategic point: when the two MP/price ratios aren't "
        "equal, the firm should reallocate."
    ))


def slide_36(prs):
    """Is Rivian's current plan optimal? (production function + MPL/MPK).

    2026-05-19 (manual round 4): slide 37 was merged into this slide and
    deleted.  The MPL / MPK derivation now lives directly under the table
    lookup.  Convention switched to "next-higher input level" (forward
    differences): MPL goes 2,000 → 2,500 at K=300; MPK goes 300 → 400 at
    L=2,000.  Axis labels bumped to 14 pt — L  (workers) rotated 270° so
    it reads bottom-up; K  (robots) moved a touch closer to the table top.

    2026-05-19 (manual round 5):
    - "Read Q at K=300, L=2,000 → Q=534" sub-bullet deleted by user.
    - "Current mix" bullet recoloured BLUE (#0070C0); a matching blue
      circle highlights the (K=300, L=2,000) cell in the table.
    - "Labor" bullet recoloured DARK GREEN (~#4F6128, accent3 + 50% lum);
      a matching green rounded rect frames the two cells used in the MPL
      interval (Q at L=2,000 and L=2,500, both at K=300).
    - "Robots" bullet recoloured DARK RED (#C00000); a matching red
      rounded rect frames the two cells used in the MPK interval
      (Q at K=300 and K=400, both at L=2,000).
    """
    BLUE_HL    = RGBColor(0x00, 0x70, 0xC0)
    DARK_GREEN = RGBColor(0x4F, 0x61, 0x28)
    DARK_RED   = RGBColor(0xC0, 0x00, 0x00)

    def draw(slide):
        bullets = [
            ("The production function at Rivian's new plant:", 0),
            ("Current mix:  300 robots, 2,000 workers", 0,
             {'color': BLUE_HL}),
            ("Let's derive the “ingredients” for the bang-for-the-buck rule", 0),
            ("Convention:  use the next-higher input level (Δ goes forward)", 1),
            ("Labor:   MPₗ  =  ΔQ / ΔL  =  (571 − 534) / (2,500 − 2,000)  =  0.07", 1,
             {'color': DARK_GREEN}),
            ("Robots:  MPₖ  =  ΔQ / ΔK  =  (617 − 534) / (400 − 300)  =  0.83", 1,
             {'color': DARK_RED}),
            ("Now:  Check if Rivian's proposed choice is optimal  (PollEV)", 0),
        ]
        _add_hierarchical_bullets(
            slide,
            left=MARGIN, top=Inches(1.85),
            width=Inches(9.0), height=Inches(4.8),
            items=bullets,
            size=24, sub_size=22, line_spacing_pts=10,
        )

        # Compact production-function table on the right.  with_axes=False
        # so the helper's default 10 pt labels don't render; the custom
        # 14 pt labels below take their place.
        tbl_left = Inches(9.55)
        tbl_top = Inches(2.30)
        col_w_label = Inches(0.72)
        col_w_data = Inches(0.55)
        tbl_h = Inches(3.70)
        _add_compact_pf_table(slide,
                               tbl_left=tbl_left, tbl_top=tbl_top,
                               with_axes=False)
        # K  (robots) label — 14 pt, centered over the data columns,
        # closer to the table top than the helper default.
        _add_text(slide,
                   tbl_left + col_w_label, tbl_top - Inches(0.25),
                   col_w_data * 4, Inches(0.25),
                   "K  (robots)",
                   size=14, italic=True, color=NAVY,
                   align=PP_ALIGN.CENTER, font="Calibri")
        # L  (workers) label — 14 pt, rotated 270° so it reads bottom-up.
        l_box = _add_text(slide,
                   tbl_left - Inches(0.83), tbl_top + tbl_h / 2 - Inches(0.135),
                   Inches(1.09), Inches(0.269),
                   "L  (workers)",
                   size=14, italic=True, color=NAVY,
                   align=PP_ALIGN.CENTER, font="Calibri")
        l_box.rotation = 270

        # Coloured overlays on specific table cells (matching bullet colours).
        n_rows = 1 + len(PF_L_VALS)  # 1 header + 12 data rows
        row_h = tbl_h / n_rows
        def cell_xywh(K, L):
            """Return (left, top, width, height) of the (K, L) table cell."""
            col_idx = PF_K_VALS.index(K)               # 0-indexed among K-vals
            row_idx = PF_L_VALS.index(L) + 1           # +1 for header row
            return (tbl_left + col_w_label + col_idx * col_w_data,
                    tbl_top + row_idx * row_h,
                    col_w_data, row_h)

        # 1. Blue circle around the (K=300, L=2,000) cell — current mix.
        cl, ct, cw, ch = cell_xywh(300, 2000)
        circle_d = Inches(0.55)
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            int(cl + cw / 2 - circle_d / 2),
            int(ct + ch / 2 - circle_d / 2),
            int(circle_d), int(circle_d),
        )
        circle.fill.background()
        circle.line.color.rgb = BLUE_HL
        circle.line.width = Pt(1.75)
        circle.shadow.inherit = False

        # 2. Dark-green rounded rect around (K=300, L=2,000) + (K=300, L=2,500)
        #    — the two cells used in the MPL forward-interval calculation.
        gl, gt, gw, _ = cell_xywh(300, 2000)
        _, gt2, _, gh2 = cell_xywh(300, 2500)
        g_bottom = gt2 + gh2
        pad = Inches(0.03)
        g_rect = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            int(gl - pad), int(gt - pad),
            int(gw + 2 * pad), int((g_bottom - gt) + 2 * pad),
        )
        g_rect.fill.background()
        g_rect.line.color.rgb = DARK_GREEN
        g_rect.line.width = Pt(1.75)
        g_rect.shadow.inherit = False
        try:
            g_rect.adjustments[0] = 0.18
        except Exception:
            pass

        # 3. Dark-red rounded rect around (K=300, L=2,000) + (K=400, L=2,000)
        #    — the two cells used in the MPK forward-interval calculation.
        rl, rt, _, rh = cell_xywh(300, 2000)
        rr_left, _, rr_w, _ = cell_xywh(400, 2000)
        r_right = rr_left + rr_w
        r_rect = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            int(rl - pad), int(rt - pad),
            int((r_right - rl) + 2 * pad), int(rh + 2 * pad),
        )
        r_rect.fill.background()
        r_rect.line.color.rgb = DARK_RED
        r_rect.line.width = Pt(1.75)
        r_rect.shadow.inherit = False
        try:
            r_rect.adjustments[0] = 0.35
        except Exception:
            pass

        # 2026-05-20 (manual round 6): the navy "Compare MP per dollar
        # across inputs at the current mix" takeaway bar at the bottom
        # was removed by hand — the bullets above already drive the
        # comparison; the bottom band was redundant.

        # 2026-05-19 (manual round 5): user added a gold POLL pill at the
        # bottom-right (matches the one introduced on slide 37 — primes
        # students that the PollEV is the very next slide).  Gold fill +
        # navy text + no leading dot + soft drop shadow + larger than the
        # default top-right pill (1.49" × 0.51", 16 pt POLL text).
        _draw_poll_pill(slide, position='bottom-right',
                         fill=GOLD, text_color=NAVY, dot_color=None,
                         width=Inches(1.49), height=Inches(0.51),
                         text_size=16, shadow=True)

    s = make_diagram_slide(
        prs, page_num=36,
        section_tag=SECTION_TAG_LR,
        title="Is Rivian's Current Plan Optimal?  (Production Function)",
        draw_diagram=draw,
    )
    # Hyperlink the "(link)" anchor in the compact-table caption to slide 11.
    _add_slide_link_in_slide(s, "link", SLIDE_IDX_PF_TABLE, prs=prs)
    _set_notes(s, (
        "The production function for Georgia in numbers.  At K = 300, "
        "L = 2,000 the table reads Q = 534.  Walk through the derivation: "
        "by convention we step UP to the next grid value (2,000 → 2,500 "
        "for labor, 300 → 400 for robots), so MPL = (571 − 534)/500 ≈ "
        "0.07 cars per worker per week and MPK = (617 − 534)/100 ≈ 0.83 "
        "cars per robot.  Don't reveal whether the mix is optimal yet — "
        "the next slide is the PollEV."
    ))



def slide_37(prs):
    """Poll: Is Rivian's input mix optimal?"""
    def draw(slide):
        _add_source_image(slide, 38, "rId4",
                           left=Inches(3.2), top=Inches(1.85),
                           height=Inches(5.1))
        _add_text(slide, MARGIN, Inches(7.0), RULE_W, Inches(0.3),
                   "Respond at PollEv.com/nvoigtlaender",
                   size=14, italic=True, color=GRAY,
                   align=PP_ALIGN.CENTER, font="Calibri")

    s = make_diagram_slide(
        prs, page_num=37,
        section_tag=SECTION_TAG_LR,
        title="Is Rivian's Input Mix Optimal?",
        draw_diagram=draw,
    )
    # 2026-05-19 (manual round 5): user removed the original top-right POLL
    # pill + leading gold dot, and added a larger gold pill at the bottom-
    # right instead (matches the one on slide 36 — same visual cue, same
    # geometry).  No leading dot.
    _draw_poll_pill(s, position='bottom-right',
                     fill=GOLD, text_color=NAVY, dot_color=None,
                     width=Inches(1.49), height=Inches(0.51),
                     text_size=16, shadow=True)
    _set_notes(s, (
        "Quick PollEv.  Looking at the numbers from the previous slide — "
        "current plan K = 300 robots, L = 2,000 workers (Q = 534), with "
        "MPL ≈ 0.07 and MPK ≈ 0.83 — is the current mix optimal?  Give "
        "them 30 seconds to think through the bang-for-the-buck ratios "
        "(MP / price).  Some will say yes, some no;  reveal in the next "
        "slide.  The point isn't the vote count, it's the active "
        "calculation."
    ))


def slide_38(prs):
    """Solution on optimal input mix."""
    def draw(slide):
        # 2026-05-19 (manual round 4): forward (next-higher) intervals,
        # consistent with the convention codified on slide 36:
        #   MPL over L ∈ [2000, 2500] at K=300: (571 − 534)/500 ≈ 0.07
        #   MPK over K ∈ [300, 400] at L=2000: (617 − 534)/100 ≈ 0.83
        # MP/$:  MPL/w  = 0.07/1200  ≈ 5.8e-5
        #        MPK/pK = 0.83/20000 ≈ 4.2e-5
        # Workers' bang-for-the-buck still wins by ~1.4x → shift toward
        # workers (same conclusion direction; smaller margin than the
        # lookback variant).
        rows = [
            ("",            "MP",            "Price",   "MP per $"),
            ("Robots  (K)", "≈ 0.83 cars",   "$20,000", "0.000042 cars / $"),
            ("Workers (L)", "≈ 0.07 cars",   "$1,200",  "0.000058 cars / $"),
        ]
        col_w = [Inches(3.0), Inches(2.5), Inches(2.5), Inches(3.0)]
        x0 = (SLIDE_W - sum(col_w)) // 2
        y0 = Inches(2.0)
        row_h = Inches(0.7)
        for r, row in enumerate(rows):
            cx = x0
            for c, val in enumerate(row):
                fill = NAVY if r == 0 else (RGBColor(0xF4, 0xF1, 0xEA) if r % 2 else WHITE)
                txt_color = WHITE if r == 0 else NAVY
                bold = r == 0 or c == 0
                shp = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    int(cx), int(y0 + r * row_h),
                    int(col_w[c]), int(row_h),
                )
                shp.fill.solid()
                shp.fill.fore_color.rgb = fill
                shp.line.color.rgb = RULE
                shp.line.width = Pt(0.5)
                shp.shadow.inherit = False
                tf = shp.text_frame
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                tf.margin_left = Inches(0.1)
                tf.margin_right = Inches(0.1)
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = val
                run.font.name = "Calibri"
                run.font.size = Pt(18)
                run.font.bold = bold
                run.font.color.rgb = txt_color
                cx += col_w[c]

        # Conclusion: stacked-fraction OMML formula + arrow + advice.
        # 2026-05-20 (manual round 6):
        #   - Operand order REVERSED so the LARGER side (workers') sits on
        #     the left:  MP_L / w  >  MP_K / p_K  (reverse=True, op=' > ').
        #   - Box shrunk to hug the formula and moved right-of-centre to
        #     make room for a leading "Thus:" label.  New position
        #     (5.29, 4.37) and size 2.75 × 1.20 (was 5.50 × 1.10 at 3.92,
        #     4.70 in round 5; 8.30 × 1.10 originally).
        _add_math_equation(
            slide,
            left=Inches(5.29), top=Inches(4.37),
            width=Inches(2.75), height=Inches(1.20),
            omml_content=_formula_bang_for_buck(op=' > ', reverse=True),
            size_pt=32, color=NAVY,
            fill=RGBColor(0xF4, 0xF1, 0xEA),
            line=NAVY,
        )
        # 2026-05-20 (manual round 6): new "Thus: " label sitting to the
        # left of the formula box (32 pt Calibri bold navy).
        _add_text(slide, Inches(4.04), Inches(4.65), Inches(1.04), Inches(0.54),
                   "Thus:", size=32, bold=True, color=NAVY, font="Calibri",
                   align=PP_ALIGN.CENTER)
        # 2026-05-20 (manual round 6): "Rivian should hire …" line nudged
        # right + down (MARGIN, 5.85) → (0.44, 6.02) to clear the now-
        # taller formula box.
        _add_text(slide, Inches(0.44), Inches(6.02), Inches(12.78), Inches(0.4),
                   "→  Rivian should hire more workers, fewer robots",
                   size=20, bold=True, color=NAVY,
                   align=PP_ALIGN.CENTER, font="Calibri")
        # 2026-05-20 (manual round 6): the gold "Equalize MP per $ →
        # reach the optimal mix" takeaway bar at the bottom was removed
        # by hand — the formula + advice line above carries the message.

    s = make_diagram_slide(
        prs, page_num=38,
        section_tag=SECTION_TAG_LR,
        title="Solution:  The Optimal Input Mix",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "Reveal: the mix isn't optimal.  Workers' bang-for-the-buck "
        "(MPL / w ≈ 5.8 × 10⁻⁵) is roughly 1.4× robots' (MPK / pK ≈ "
        "4.2 × 10⁻⁵), so Rivian should hire more workers and use fewer "
        "robots.  Intuition: at K = 300, L = 2,000 the marginal worker "
        "delivers more output per dollar than the marginal robot — "
        "robots are relatively over-priced for the current labor force. "
        "Same bang-for-the-buck rule as before; the conclusion direction "
        "depends on where the operating point sits."
    ))


def slide_39(prs):
    """When prices change, the input mix shifts: robot tax & union wages.

    The two body boxes use OMML for the math (p_K, MP_K / p_K with real
    subscripts + stacked fractions) and Calibri for the descriptive
    'shift toward …' line beneath each formula.
    """
    def draw(slide):
        _add_text(slide, MARGIN, Inches(1.9), RULE_W, Inches(0.55),
                   "What's the effect on the optimal input mix if…",
                   size=24, italic=True, color=GRAY,
                   align=PP_ALIGN.CENTER, font="Calibri")

        col_w = Inches(6.0)
        col_h = Inches(3.4)
        gap = Inches(0.4)
        left_x = (SLIDE_W - 2 * col_w - gap) // 2
        right_x = left_x + col_w + gap
        y = Inches(2.7)
        cream = RGBColor(0xF4, 0xF1, 0xEA)
        body_h = col_h - Inches(0.75)

        def _column(x, header_text, eq_xml, conclusion, math_pos=None):
            # Header band
            _add_filled_box(slide, x, y, col_w, Inches(0.7),
                             header_text,
                             fill=NAVY, text_color=WHITE,
                             size=18, bold=True)
            # Cream body background (no text)
            _add_filled_box(slide, x, y + Inches(0.75),
                             col_w, body_h, "",
                             fill=cream, text_color=NAVY,
                             size=20, bold=False)
            # OMML equation — defaults to filling the cream body's upper
            # half; ``math_pos=(left, top, width, height)`` lets callers
            # nudge the equation tighter and inset (slide 39 round 7).
            if math_pos is None:
                m_left, m_top, m_w, m_h = (
                    x, y + Inches(0.95), col_w, Inches(1.45),
                )
            else:
                m_left, m_top, m_w, m_h = math_pos
            _add_math_equation(
                slide, m_left, m_top, m_w, m_h,
                eq_xml, size_pt=24, color=NAVY,
            )
            # Conclusion line beneath the formula
            _add_text(slide, x, y + Inches(2.55),
                       col_w, Inches(0.5),
                       conclusion,
                       size=20, color=NAVY, font="Calibri",
                       align=PP_ALIGN.CENTER)

        # 2026-05-20 (manual round 7): user hand-tightened the math
        # boxes — both shrunk and inset inside their cream bodies, with
        # non-breaking-space tweaks to the operator spacing.  Left
        # formula slightly wider with 3 nbsp + trailing nbsp; right
        # formula a touch narrower with 2 nbsp and no trailing nbsp.

        # LEFT column: tax on robots → p_K rises → MP_K / p_K falls
        eq_left = (
            _omml_sub(_omml_run('p'), _omml_run('K')) +
            _omml_text('↑\xa0\xa0\xa0⇒\xa0') +
            _omml_frac(
                _omml_sub(_omml_run('MP'), _omml_run('K')),
                _omml_sub(_omml_run('p'), _omml_run('K'))
            ) +
            _omml_text('\xa0↓')
        )
        _column(
            left_x,
            "The government introduces a high tax on robots",
            eq_left,
            "→  shift toward more labor",
            math_pos=(Inches(2.17), Inches(4.00),
                      Inches(3.00), Inches(0.93)),
        )

        # RIGHT column: union wages → w rises → MP_L / w falls
        eq_right = (
            _omml_run('w') +
            _omml_text('↑\xa0\xa0⇒\xa0') +
            _omml_frac(
                _omml_sub(_omml_run('MP'), _omml_run('L')),
                _omml_run('w')
            ) +
            _omml_text('↓')
        )
        _column(
            right_x,
            "Labor unions demand significantly higher wages",
            eq_right,
            "→  shift toward more automation",
            math_pos=(Inches(8.73), Inches(4.15),
                      Inches(2.55), Inches(0.86)),
        )

        # Bottom takeaway — 2026-05-20 (manual round 7): user rewrote
        # the line from "…the cheaper input" → "…the input that has
        # become cheaper" for clearer phrasing.
        _add_takeaway_bar(slide,
                           "When input prices change, the optimal mix shifts toward the input that has become cheaper",
                           top=Inches(6.38), fill=GOLD, text_color=NAVY,
                           width=Inches(12.0), size=18,
                           rounded=True, shadow=True)

    s = make_diagram_slide(
        prs, page_num=39,
        section_tag=SECTION_TAG_LR,
        title="When Prices Change, the Input Mix Shifts",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "Comparative statics. When input prices change, the optimal mix "
        "shifts: a tax on robots pushes Rivian toward more labor; rising "
        "wages push them toward more automation. Real strategic "
        "implications for any firm facing input-price shocks – including AI "
        "labs deciding between GPU spend and engineer headcount."
    ))


def slide_40(prs):
    """'Bang for the Buck' in grocery shopping.

    2026-05-21: rebuilt to mirror original slide 47 (same wording, same
    OMML formulas, same image sizes; pictures shifted right by ~1.67"
    to centre the pair on the 16:9 deck while preserving the original's
    relative spacing).  The bottom NAVY 'universal decision rule' bar
    from the prior rebuild is kept and upgraded with rounded corners +
    a drop shadow.
    """
    def draw(slide):
        # ---- Bullets (original 47 wording, verbatim) ----
        bullets = [
            ("P&G offers two sizes for its flagship detergent", 0),
            ("Small Tide Ultra Oxi (29 loads) for $17", 1),
            ("Large Tide Ultra Oxi (81 loads) for $35", 1),
        ]
        _add_hierarchical_bullets(
            slide,
            left=MARGIN, top=Inches(1.55),
            width=Inches(12.0), height=Inches(1.1),
            items=bullets,
            size=24, sub_size=22, line_spacing_pts=8,
        )

        # ---- Three OMML formulas, stacked, full-width centered ----
        # MP_x / p_x  — multi-letter "small" / "large" subscripts kept
        # italic (math-default) to match the original 47 rendering.
        def mp_over_p(label):
            return _omml_frac(
                _omml_sub(_omml_run('MP'), _omml_run(label)),
                _omml_sub(_omml_run('p'),  _omml_run(label)),
            )

        f_small_eq = (
            mp_over_p('small')
            + _omml_text(' = ')
            + _omml_frac(_omml_text('29'), _omml_text('17'))
            + _omml_text(' = ')
            + _omml_text('1.7')
        )
        f_large_eq = (
            mp_over_p('large')
            + _omml_text(' = ')
            + _omml_frac(_omml_text('81'), _omml_text('35'))
            + _omml_text(' = ')
            + _omml_text('2.3')
        )
        f_compare = (
            mp_over_p('small')
            + _omml_text(' &lt; ')      # bare '<' would break the XML parser
            + mp_over_p('large')
        )

        # ---- Three OMML formulas — positions/sizes per 2026-05-21
        # hand-tweaks in PowerPoint (small ratio left, large ratio right,
        # inequality centred below the pair). ----
        # F1 small ratio:  was full-width row at (0, 2.70, SLIDE_W, 0.50)
        _add_math_equation(
            slide, left=Inches(1.720), top=Inches(3.061),
            width=Inches(3.130), height=Inches(0.859),
            omml_content=f_small_eq, size_pt=22, color=NAVY,
        )
        # F2 large ratio:  was full-width row at (0, 3.25, SLIDE_W, 0.50)
        _add_math_equation(
            slide, left=Inches(8.080), top=Inches(3.005),
            width=Inches(3.363), height=Inches(0.910),
            omml_content=f_large_eq, size_pt=22, color=NAVY,
        )
        # F3 inequality:  was full-width row at (0, 3.80, SLIDE_W, 0.50)
        _add_math_equation(
            slide, left=Inches(4.780), top=Inches(4.450),
            width=Inches(3.300), height=Inches(0.910),
            omml_content=f_compare, size_pt=22, color=NAVY,
        )

        # ---- Two product pictures — positions per 2026-05-21 hand-tweaks
        # (small bottle nudged left/up; large bottle nudged up). ----
        # Source-deck rels: rId5 = small bottle, rId6 = large bottle.
        # Small:  was (3.58, 4.39)  →  (2.69, 4.165)
        # Large:  was (8.69, 4.26)  →  (8.69, 4.095)
        _add_source_image(slide, 41, "rId5",
                           left=Inches(2.690), top=Inches(4.165),
                           height=Inches(1.90))
        _add_source_image(slide, 41, "rId6",
                           left=Inches(8.690), top=Inches(4.095),
                           height=Inches(2.16))

        # ---- Bottom NAVY takeaway bar — kept from prior rebuild,
        # now rounded + drop-shadowed (2026-05-21 user ask). ----
        _add_takeaway_bar(slide,
                           "Bang-for-the-buck:  a universal decision rule",
                           top=Inches(6.5), fill=NAVY, width=Inches(9.0),
                           rounded=True, shadow=True)

    s = make_diagram_slide(
        prs, page_num=40,
        section_tag=SECTION_TAG_LR,
        title="“Bang for the Buck” in Grocery Shopping",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "The bang-for-the-buck rule isn't just for factories – you apply "
        "it every week at the grocery store. P&G offers two sizes of Tide "
        "Ultra Oxi: small at 29 loads for $17, large at 81 loads for $35. "
        "Compute loads-per-dollar: 29/17 ≈ 1.7 for the small, 81/35 "
        "≈ 2.3 for the large. The large size delivers more cleaning "
        "per dollar, so the bang-for-the-buck rule says to buy the large "
        "bottle. Same logic as the factory's input choice, just a "
        "different decision."
    ))


def slide_41(prs):
    """Section divider – Part 2: Costs."""
    s = make_section_agenda(
        prs, page_num=41,
        current_part_idx=1,        # Part 2 now active
        section_tag=SECTION_TAG_DIV,
        title="Part 2:  Costs – Producing at the Lowest Price",
    )
    _set_notes(s, (
        "Part 2 – Costs. We've covered what to PRODUCE; now we cover what "
        "it COSTS, and crucially, which costs actually matter for "
        "decisions."
    ))


# --------------------------------------------------------------------------
# §2.1 Cost Concepts (slides 43-62)
# --------------------------------------------------------------------------

SECTION_TAG_P2 = "Module 3 · Costs · Cost Concepts"
SECTION_TAG_P2_LR = "Module 3 · Costs · Long-Run & Scale"


def slide_42(prs):
    """Cost types: variable / fixed / sunk, each with examples; takeaway
    'sunk costs should not affect managerial decisions'."""
    def draw(slide):
        # 2026-05-21: distinct accents per cost type — DARK_GREEN on
        # Variable (always-consider), NAVY on Fixed (conditional), DARK_RED
        # on Sunk (always-ignore). 2026-05-21 v2: re-ordered left→right by
        # decision weight (Variable → Fixed → Sunk) and added an
        # "Examples" card under each column.
        DARK_RED   = RGBColor(0x8B, 0x1A, 0x1A)
        DARK_GREEN = RGBColor(0x2E, 0x7D, 0x32)

        # Each band carries (label, body_paragraphs, examples, fill).
        # body_paragraphs is a list of (text, space_before_pt, space_after_pt)
        # — first entry replaces the default paragraph created by
        # _add_outlined_box; the rest are appended via add_paragraph.
        # 2026-05-21 v3 hand-edits ported from PowerPoint:
        #   • Variable verdict: "Always consider!" (was just "Consider!").
        #   • Fixed verdict: split from one soft-broken paragraph into two
        #     real paragraphs — "Consider for long-run decisions" then
        #     "(entry, exit, capacity)" without the soft break.
        #   • Sunk verdict: 12 pt moved from space-before → space-after, so
        #     the centered block shifts up inside the box.
        bands = [
            ("Variable Costs",
             [("Depend on volume produced (Q)", None, None),
              ("→  Always consider!", 12, None)],
             ["Buy raw materials", "Hire workers",
              "Shipping & packaging"],
             DARK_GREEN),
            ("Fixed Costs",
             [("Do not depend on quantity produced (Q)", None, None),
              ("→  Consider for long-run decisions", 12, None),
              ("(entry, exit, capacity)", None, None)],
             ["Rent / lease payments", "Insurance & property tax",
              "R&D investments (for future innovation)"],
             NAVY),
            ("Sunk Costs",
             [("A fixed cost that ", None, None),
              ("cannot be recovered\n→  Ignore!", None, 12)],
             ["Past R&D", "Past advertising",
              "Non-refundable deposits"],
             DARK_RED),
        ]
        band_w = Inches(3.95)
        gap = Inches(0.15)
        total_w = band_w * 3 + gap * 2
        start_x = (SLIDE_W - total_w) // 2

        for i, (label, body_paras, examples, fill) in enumerate(bands):
            bx = start_x + (band_w + gap) * i
            # Header band
            hdr = _add_filled_box(
                slide, bx, Inches(2.05), band_w, Inches(0.7), label,
                fill=fill, text_color=WHITE, size=22, bold=True,
            )
            _add_drop_shadow(hdr)
            # Body description — first paragraph supplied to the helper,
            # remaining verdict paragraphs appended below.
            first_text, first_sb, first_sa = body_paras[0]
            bdy = _add_outlined_box(
                slide, bx, Inches(2.75), band_w, Inches(1.5), first_text,
                fill=WHITE, line=fill, text_color=NAVY,
                size=18, bold=False, line_w=1.5,
            )
            _add_drop_shadow(bdy)
            p0 = bdy.text_frame.paragraphs[0]
            if first_sb is not None: p0.space_before = Pt(first_sb)
            if first_sa is not None: p0.space_after = Pt(first_sa)
            for text, sb, sa in body_paras[1:]:
                pn = bdy.text_frame.add_paragraph()
                pn.alignment = PP_ALIGN.CENTER
                if sb is not None: pn.space_before = Pt(sb)
                if sa is not None: pn.space_after = Pt(sa)
                rn = pn.add_run()
                rn.text = text
                rn.font.name = "Calibri"
                rn.font.size = Pt(18)
                rn.font.bold = False
                rn.font.color.rgb = NAVY

            # Examples card directly below — rounded white rect, navy
            # border, soft drop shadow. "Examples" label sits top-left;
            # 2–3 bullets follow.
            ex = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                int(bx), int(Inches(4.35)),
                int(band_w), int(Inches(1.20)),
            )
            try: ex.adjustments[0] = 0.10
            except Exception: pass
            ex.fill.solid()
            ex.fill.fore_color.rgb = WHITE
            ex.line.color.rgb = NAVY
            ex.line.width = Pt(1.25)
            ex.shadow.inherit = False
            _add_drop_shadow(ex)
            etf = ex.text_frame
            etf.word_wrap = True
            etf.margin_left = Inches(0.13)
            etf.margin_right = Inches(0.08)
            etf.margin_top = Inches(0.06)
            etf.margin_bottom = Inches(0.05)
            etf.vertical_anchor = MSO_ANCHOR.TOP
            lbl_p = etf.paragraphs[0]
            lbl_p.alignment = PP_ALIGN.LEFT
            lbl_r = lbl_p.add_run()
            lbl_r.text = "Examples"
            lbl_r.font.name = "Calibri"
            lbl_r.font.size = Pt(12)
            lbl_r.font.bold = True
            lbl_r.font.italic = True
            lbl_r.font.color.rgb = NAVY
            for ex_text in examples:
                pe = etf.add_paragraph()
                pe.alignment = PP_ALIGN.LEFT
                re_ = pe.add_run()
                re_.text = "•  " + ex_text
                re_.font.name = "Calibri"
                re_.font.size = Pt(14)
                re_.font.bold = False
                re_.font.color.rgb = NAVY

        # 2026-05-21 v3 hand-edit: takeaway moved to sit directly under
        # the Sunk Costs column (right edge), narrowed to band_w so it
        # visually anchors to that single column; text wraps to 2 lines
        # (hence the taller 0.82 in height). T=5.91 in matches the user's
        # PowerPoint drag.
        sunk_left = start_x + (band_w + gap) * 2
        _add_takeaway_bar(
            slide,
            "Sunk costs should not affect managerial decisions",
            left=sunk_left, top=5405247,
            width=band_w, height=747521,
            fill=DARK_RED, text_color=WHITE,
            rounded=True, shadow=True,
        )

    s = make_diagram_slide(
        prs, page_num=42,
        section_tag=SECTION_TAG_P2,
        title="Three Cost Types",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "The single most important cost concept for executives, in five "
        "words: sunk costs are not costs. They've already been spent; they "
        "cannot be recovered. Any forward-looking decision should ignore "
        "them. Period."
    ))


def slide_43(prs):
    """Group work: choosing a car (own car vs. company car)."""
    def draw(slide):
        # Setup paragraph
        _add_text(slide, MARGIN, Inches(1.85), RULE_W, Inches(0.55),
                  "For a business trip, you can use:",
                  size=22, bold=True, color=NAVY, font="Calibri")

        # Two options side by side
        opt_w = Inches(5.5)
        opt_h = Inches(1.4)
        opt_gap = Inches(0.6)
        start_x = (SLIDE_W - opt_w * 2 - opt_gap) // 2
        own_left = start_x
        co_left = start_x + opt_w + opt_gap
        opt_top = Inches(2.45)

        # 2026-05-21 hand-edits on the option boxes:
        #   • Title bumped to 28 pt (was 20 pt).
        #   • Subtitle wording rewritten — "you are reimbursed 50¢ / mile"
        #     (was "+  reimbursed 50¢ / mile") and "(full cost incl.
        #     charge paid by your company)" (was "…gas covered)").
        #   • A leading newline opens visible vertical air between title
        #     and subtitle inside the box.
        def _draw_option_box(left, title, subtitle):
            shp = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, int(left), int(opt_top),
                int(opt_w), int(opt_h),
            )
            shp.fill.solid()
            shp.fill.fore_color.rgb = WHITE
            shp.line.color.rgb = NAVY
            shp.line.width = Pt(1.5)
            shp.shadow.inherit = False
            tf = shp.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.1)
            tf.margin_right = Inches(0.1)
            tf.margin_top = Inches(0.05)
            tf.margin_bottom = Inches(0.05)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p0 = tf.paragraphs[0]
            p0.alignment = PP_ALIGN.CENTER
            r0 = p0.add_run()
            r0.text = title
            r0.font.name = "Calibri"
            r0.font.size = Pt(28)
            r0.font.bold = True
            r0.font.color.rgb = NAVY
            p1 = tf.add_paragraph()
            p1.alignment = PP_ALIGN.CENTER
            r1 = p1.add_run()
            r1.text = "\n" + subtitle
            r1.font.name = "Calibri"
            r1.font.size = Pt(20)
            r1.font.bold = True
            r1.font.color.rgb = NAVY

        _draw_option_box(own_left, "Your own car",
                         "→  you are reimbursed 50¢ / mile")
        _draw_option_box(co_left, "Company car",
                         "(full cost incl. charge paid by your company)")

        # 2026-05-21 hand-edits to the wrapper container:
        #   • Top nudged down from 3.85 → 4.00 so a visible gap opens
        #     between the "Your own car" box (bottom 3.85) and the
        #     wrapper.
        #   • Border thickened from 1.0 pt → 2.0 pt for more presence.
        #   • Header text bumped from 14 pt → 18 pt.
        #   • Cost boxes shrunk in height (0.75 → 0.40 in, ≈50%).
        #   • 2026-05-21 v2: wrapper height shrunk 1.95 → 1.50 in; cost
        #     rows pulled up to sit just below the header (instead of
        #     pinned to the wrapper bottom); the "Should you use…"
        #     question pushed further down relative to the now-shorter
        #     wrapper.
        wrap_l = own_left
        wrap_t = Inches(4.00)
        wrap_w = opt_w
        wrap_h = Inches(1.50)  # bottom at Y = 5.50
        wrap = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            int(wrap_l), int(wrap_t), int(wrap_w), int(wrap_h),
        )
        try: wrap.adjustments[0] = 0.04
        except Exception: pass
        wrap.fill.solid()
        wrap.fill.fore_color.rgb = WHITE
        wrap.line.color.rgb = NAVY
        wrap.line.width = Pt(2.0)
        wrap.shadow.inherit = False

        # Header sits at the top of the wrapper, 18 pt italic gray
        _add_text(slide, wrap_l, Inches(4.10), wrap_w, Inches(0.40),
                  "Costs associated with your car  (per mile driven):",
                  size=18, italic=True, color=GRAY, font="Calibri",
                  align=PP_ALIGN.CENTER)

        # Four cost boxes — narrow (W=2.50 in, hugs longest text
        # "45¢   lease on the vehicle") and short (H=0.40 in, ≈50% of
        # the prior 0.75 in). 2026-05-21 v2: rows pulled up so row 1 sits
        # flush below the header (gap ~0.02 in), giving the wrapper a
        # tight header-then-grid stack with breathing room below.
        costs = [
            ("20¢", "insurance"),
            ("20¢", "maintenance"),
            ("15¢", "electricity"),
            ("45¢", "lease on the vehicle"),
        ]
        cost_w = Inches(2.50)
        cost_h = Inches(0.40)
        col_gap = Inches(0.10)
        row_gap = Inches(0.05)
        grid_w = cost_w * 2 + col_gap
        grid_x0 = wrap_l + (wrap_w - grid_w) // 2
        row1_top = Inches(4.5225)
        for i, (amt, lbl) in enumerate(costs):
            row, col = divmod(i, 2)
            cx = grid_x0 + col * (cost_w + col_gap)
            cy = row1_top + row * (cost_h + row_gap)
            _add_filled_box(
                slide, cx, cy, cost_w, cost_h,
                f"{amt}   {lbl}",
                fill=NAVY, text_color=WHITE,
                size=15, bold=True,
            )

        # Question — prefixed with double-right-arrow (⇒). 2026-05-21 v2
        # hand-edit: T 5.95 → 5.815 in (pushed further down RELATIVE to
        # the now-shorter wrapper, which now ends at 5.50).
        _add_text(slide, MARGIN, Inches(5.815), RULE_W, Inches(0.45),
                  "⇒  Should you use your own car or the company car?",
                  size=22, bold=True, color=NAVY, font="Calibri",
                  align=PP_ALIGN.CENTER)
        # Discussion-break badge — 2026-05-21 v2 hand-edit: T = 6.34 in
        # (was Inches(6.25) + Cm(0.4) = 6.407 in; user nudged up), and
        # shifted right ~0.20 in (left margin from slide right edge
        # tightened from MARGIN ≈ 0.28 in to ~0.08 in).
        _add_discussion_break(slide, top=Inches(6.34), width=Inches(4.8),
                              left=Inches(8.457))

    s = make_diagram_slide(
        prs, page_num=43,
        section_tag=SECTION_TAG_P2,
        title="Group Work:  Your Car or the Company Car?",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "A group-work exercise. Two cars, same daily operating cost, but "
        "different sunk amounts. Which should you drive today? Whatever "
        "your gut says, the answer is: ignore the sunk cost – it's the "
        "same either way (the lease is sunk; only the variable per-mile "
        "costs matter for the marginal trip decision)."
    ))


def slide_44(prs):
    """Why studios finish movies they know lose money — Waterworld (1995).

    2026-05-21 hand-edits ported from PowerPoint:
      • Title reworded.
      • "Waterworld (1995)" header textbox added above the picture
        (28 pt bold navy, centered).
      • Picture moved left/down (L 3.70 → 2.65 in, T 1.65 → 2.12 in).
      • Bottom caption removed.
    """
    def draw(slide):
        # "Waterworld (1995)" header above the picture
        _add_text(slide, Inches(3.631), Inches(1.554),
                  Inches(6.665), Inches(0.572),
                  "Waterworld (1995)",
                  size=28, bold=True, color=NAVY, font="Calibri",
                  align=PP_ALIGN.CENTER)
        # The Waterworld poster, large and re-positioned
        _add_source_image(slide, 45, "rId3",
                          left=Inches(2.65), top=Inches(2.12),
                          height=Inches(4.85))

    s = make_diagram_slide(
        prs, page_num=44,
        section_tag=SECTION_TAG_P2,
        title="Why Studios Finish Movies They Know Lose Money",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "Hollywood case. Kevin Costner's Waterworld – they knew it would "
        "flop, why release it anyway? Answer: even a flop adds revenue "
        "net of marketing/release costs. The hundreds of millions already "
        "spent on production are sunk. Decision should be forward-looking "
        "only."
    ))


def slide_45(prs):
    """Sunk cost in Waterworld – decision tree across 3 scenarios.

    2026-05-21 v2 hand-edits ported:
      • Title reworded.
      • Whole table shifted up ~0.28 in (row_y starts at 1.569).
      • "Expected Additional Cost" data cells get a soft light-yellow
        fill — visually flags the row whose value (<150 = revenue) is
        the pedagogical hinge of "make the film".
      • "Make the film!" decision cells get a NAVY outline (in addition
        to the gold fill) so they read as part of the same boxed grid.
      • Takeaway: reworded "Ignore sunk costs — continue whenever
        Expected Revenue > Expected Additional Cost"; nudged up to
        T=6.389 in; now rounded with a soft drop shadow.
    """
    def draw(slide):
        # Soft light-yellow for the "Expected Additional Cost" row —
        # noticeably lighter than the GOLD on the decision cells.
        LIGHT_YELLOW = RGBColor(0xFF, 0xF2, 0xCC)

        # Each scenario: (header, revenue, sunk, additional, overall, profit, decision)
        scenarios = [
            ("June 1994",       "150", "16",  "84    (<150)", "100", "+50", "Make the film!"),
            ("September 1994",  "150", "100", "40    (<150)", "140", "+10", "Make the film!"),
            ("December 1994",   "150", "140", "35    (<150)", "175", "−25", "Make the film!"),
        ]
        # Original-style row labels (in the column-0 cells of the table).
        row_labels = [
            "",                                  # corner cell
            "Expected Revenues  ($M)",
            "Sunk Cost  ($M)",
            "Expected Additional Cost  ($M)",
            "Overall Cost (incl. sunk)  ($M)",
            "Expected Profit  ($M)",
            "Decision",
        ]
        # 4-column layout: wide label col + 3 narrower scenario cols.
        label_w = Inches(4.00)
        col_w   = Inches(2.50)
        col_gap = Inches(0.10)
        total_w = label_w + col_gap + col_w * 3 + col_gap * 2
        table_l = (SLIDE_W - total_w) // 2
        col_x0  = table_l + label_w + col_gap

        # Vertical row layout — 2026-05-21 v2: whole table shifted up
        # ~0.28 in from the prior version.
        row_y = [Inches(1.569), Inches(2.119), Inches(2.819),
                  Inches(3.519), Inches(4.219), Inches(4.919),
                  Inches(5.619)]
        row_h = [Inches(0.50), Inches(0.60), Inches(0.60),
                  Inches(0.60), Inches(0.60), Inches(0.60),
                  Inches(0.55)]

        # Left column: row labels in NAVY-outlined cells, left-aligned.
        for i, lbl in enumerate(row_labels):
            if i == 0:
                # Empty corner cell (white fill, navy border)
                _add_outlined_box(
                    slide, table_l, row_y[0], label_w, row_h[0],
                    "", fill=WHITE, line=NAVY, text_color=NAVY,
                    size=10, bold=False, line_w=1.0,
                )
                continue
            is_decision = (i == 6)
            shp = _add_outlined_box(
                slide, table_l, row_y[i], label_w, row_h[i],
                lbl, fill=WHITE, line=NAVY, text_color=NAVY,
                size=18 if not is_decision else 20,
                bold=is_decision, line_w=1.0,
            )
            tf = shp.text_frame
            tf.paragraphs[0].alignment = PP_ALIGN.LEFT
            tf.margin_left = Inches(0.18)

        # 3 scenario columns
        for j, sc in enumerate(scenarios):
            x = col_x0 + (col_w + col_gap) * j
            # Column header (date band)
            _add_filled_box(
                slide, x, row_y[0], col_w, row_h[0],
                sc[0], fill=NAVY, text_color=WHITE,
                size=20, bold=True,
            )
            # 5 data rows: revenue, sunk, additional, overall, profit.
            # i==2 is the "Expected Additional Cost" row — flagged with
            # a soft light-yellow fill (the pedagogical "this matters"
            # row: additional cost is below revenue 150 in every case).
            for i, v in enumerate(sc[1:6]):
                cell_fill = LIGHT_YELLOW if i == 2 else WHITE
                _add_outlined_box(
                    slide, x, row_y[i + 1], col_w, row_h[i + 1],
                    v, fill=cell_fill, line=NAVY, text_color=NAVY,
                    size=18, bold=False, line_w=1.0,
                )
            # Decision band (gold fill + navy outline so it reads as
            # part of the same boxed grid).
            _add_filled_box(
                slide, x, row_y[6], col_w, row_h[6],
                sc[6], fill=GOLD, text_color=NAVY, line=NAVY,
                size=18, bold=True,
            )

        # Bottom takeaway — rounded with soft drop shadow.
        _add_takeaway_bar(
            slide,
            "Ignore sunk costs —  continue whenever          "
            "Expected Revenue  >  Expected Additional Cost",
            top=Inches(6.389), fill=NAVY, text_color=WHITE,
            width=Inches(11.5), size=18,
            rounded=True, shadow=True,
        )

    s = make_diagram_slide(
        prs, page_num=45,
        section_tag=SECTION_TAG_P2,
        title="Waterworld:  Sunk Costs and Decisions over Time",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "The decision tree behind the Waterworld decision. Three scenarios "
        "– optimistic, neutral, pessimistic – differ in how much was sunk "
        "vs. how much more it costs to release. In ALL three the "
        "release-anyway revenue exceeds the marginal release cost. So "
        "they released. Sunk costs are sunk."
    ))


def slide_46(prs):
    """Modern sunk cost: Meta's Reality Labs has lost $70B+.

    Layout: bullets top-left, WSJ headline screenshot bottom-left,
    Zuckerberg-with-Quest photo right, rounded takeaway pill bottom-right.

    Reworked 2026-05-22 to reflect the user's hand-edits in PowerPoint:
      • Title shortened — "$50B+ Since 2020" → "$70B+" (matches WSJ
        screenshot, which cites $77B losses).
      • Bullets reworked to lead with the spend, then explain *why*
        Reality Labs is effectively sunk (tech not commercially
        viable, weak adoption, payoff slipping), ending with Meta's
        2026 decision to shift investment to AI.
      • Old Meta Quest 3 source-deck image replaced by a Zuckerberg /
        Quest 3 photo (now rounded corners + drop shadow).
      • WSJ "Meta Plans to Shift Spending Away from the Metaverse"
        headline added bottom-left, with drop shadow + thin gray
        border (kept rectangular per user's hand-positioning).
      • Old full-width takeaway bar replaced by a smaller rounded
        navy/gold pill in the lower-right corner.
    """
    def draw(slide):
        # ---- Bullets (top-left) ----------------------------------------
        # User reworked the first major bullet + its two sub-points to
        # focus on commercial viability rather than the original "Wall
        # Street vs. Zuckerberg" framing.  Last bullet ("In 2026 …") was
        # hand-added by the user.
        bullets = [
            ("Meta poured ~$70B into Reality Labs since 2020", 0),
            ("Quest headsets, Horizon Worlds – thin adoption, little revenue", 1),
            ("Hardware/software still not commercially viable – payoff keeps slipping", 1),
            ("Wall Street kept asking when it pays off", 0),
            ("In 2026: Decision to shift spending to AI", 0),
        ]
        # 2026-05-22 hand-tweaked from H 4.9 → 3.25 in (originally 3.52,
        # tightened further on second pass) to free up the lower-left
        # quadrant for the WSJ screenshot.  sub_size hand-bumped from
        # 22 → 24 pt — user wanted the two viability sub-bullets to
        # read at the same weight as the main bullets for EMBA legibility.
        _add_hierarchical_bullets(
            slide,
            left=MARGIN, top=Inches(1.57),
            width=Inches(7.0), height=Inches(3.25),
            items=bullets,
            size=24, sub_size=24, line_spacing_pts=12,
        )

        # ---- Zuckerberg / Reality Labs photo (right) -------------------
        # 2026-05-22 hand-added by user — Zuckerberg holding/using a
        # Quest-class headset.  Rounded corners + soft drop shadow so
        # the photo "lifts" off the slide (per course CLAUDE.md).
        zuck_path = OUT_DIR / "_zuckerberg_realitylabs.png"
        if zuck_path.exists():
            zuck = slide.shapes.add_picture(
                str(zuck_path),
                int(Inches(7.37)), int(Inches(2.16)),
                width=int(Inches(5.77)), height=int(Inches(3.76)),
            )
            _apply_picture_style(zuck, corner_pct=8)

        # ---- WSJ headline screenshot (bottom-left) ---------------------
        # 2026-05-22 hand-added by user — screenshot of WSJ Tech/AI
        # article "Meta Plans to Shift Spending Away from the Metaverse"
        # citing $77B in Reality Labs losses.  Position preserved
        # exactly as user placed it.  Drop shadow + thin gray border
        # per user request (no rounded corners — it's a doc screenshot).
        wsj_path = OUT_DIR / "_wsj_meta_reality_labs.png"
        if wsj_path.exists():
            wsj = slide.shapes.add_picture(
                str(wsj_path),
                int(Inches(0.28)), int(Inches(5.23)),
                width=int(Inches(6.57)), height=int(Inches(1.78)),
            )
            _add_drop_shadow(wsj)
            wsj.line.color.rgb = GRAY
            wsj.line.width = Pt(0.75)

        # ---- Takeaway pill (lower-right) -------------------------------
        # 2026-05-22 hand-positioned by user to the lower-right corner
        # (was: full-width takeaway bar at bottom).  Rounded corners +
        # soft drop shadow per user request.
        _add_rounded_filled_box(
            slide,
            Inches(7.28), Inches(6.40),
            Inches(5.89), Inches(0.55),
            label="→ Classic sunk-cost discipline, 2020s edition",
            fill=GOLD, text_color=NAVY,
            size=18, bold=True,
            corner_pct=0.20, shadow=True,
        )

    s = make_diagram_slide(
        prs, page_num=46,
        section_tag=SECTION_TAG_P2,
        title="Modern Sunk Cost:  Meta's Reality Labs Has Lost $70B+",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "The same Waterworld logic, in a current strategic context. "
        "Meta has poured roughly $70B – the Wall Street Journal puts "
        "the figure north of $77B – into Reality Labs since 2020. "
        "Quest headsets and Horizon Worlds never crossed into mass "
        "adoption, and the AR/VR hardware-and-software stack still "
        "is not commercially viable. Wall Street kept asking when it "
        "pays off, and forward-looking returns kept getting pushed "
        "out. In 2026 Meta made the call to shift incremental "
        "spending toward AI – higher MPL on the next dollar, "
        "regardless of how many billions had already been spent on "
        "Reality Labs. Classic sunk-cost discipline, 2020s edition."
    ))


def slide_47(prs):
    """Sunk cost & opportunity cost: Apple's canceled Apple Car.

    User hand-edits ported 2026-05-22:
      • Title broadened — was "Opportunity Cost Is a Real Cost: …",
        now leads with BOTH cost concepts since the slide makes the
        sunk-vs-opportunity point together.
      • Bullet 2: "Sunk costs ≠ a reason to keep going" →
        "Sunk costs not a reason to keep going" (avoid the math
        glyph in a plain-prose bullet).
      • Bullet 4: "(the higher-MPL use)" → "(with higher expected
        returns)" — plainer language for EMBA pacing.
      • Bullets box: H 4.5 → 3.68 in; font 24/22 → 28/24 pt
        (heavier read for the executive audience).
      • Bottom takeaway: rewritten to "next-best alternative use
        for the same dollars & engineers", repositioned to
        L 1.53 / T 6.01 / W 10.5, and upgraded to rounded corners
        + soft drop shadow per the deck-wide takeaway treatment.
    """
    def draw(slide):
        bullets = [
            ("Apple killed Project Titan in 2024 after ~10 years and ~$10B spent", 0),
            ("Sunk costs not a reason to keep going", 0),
            ("Real reason to stop: opportunity cost of capital + 2,000 engineers", 0),
            ("Reallocated → AI / Apple Intelligence (with higher expected returns)", 1),
        ]
        # 2026-05-22 hand-tweaks: H 4.5 → 3.68 in; size 24/22 → 28/24 pt.
        _add_hierarchical_bullets(
            slide,
            left=MARGIN, top=Inches(1.85),
            width=Inches(7.0), height=Inches(3.68),
            items=bullets,
            size=28, sub_size=24, line_spacing_pts=12,
        )

        # Vanarama Apple Car concept render on the right
        _add_source_image(slide, 48, "rId3",
                          left=Inches(7.45), top=Inches(2.2),
                          width=Inches(5.55))
        _add_text(slide, Inches(7.45), Inches(5.25), Inches(5.55), Inches(0.35),
                  "Vanarama Apple Car concept  (fair use, © Vanarama)",
                  size=11, italic=True, color=GRAY, font="Calibri",
                  align=PP_ALIGN.CENTER)

        # 2026-05-22 hand-tweaked: was a flat full-width takeaway bar
        # at top=Inches(6.5); user wanted the deck-standard rounded +
        # shadow treatment, repositioned slightly up and right.
        _add_rounded_filled_box(
            slide,
            Inches(1.53), Inches(6.01),
            Inches(10.50), Inches(0.55),
            label="Opportunity cost  =  the next-best alternative use for the same dollars  &  engineers",
            fill=NAVY, text_color=WHITE,
            size=20, bold=True,
            corner_pct=0.20, shadow=True,
        )

    s = make_diagram_slide(
        prs, page_num=47,
        section_tag=SECTION_TAG_P2,
        title="Sunk Cost & Opportunity Cost:  Apple's Canceled Apple Car",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "Apple killed Project Titan in 2024 after roughly a decade and "
        "$10B spent on the car. The slide makes two cost points at "
        "once. First, the sunk costs are sunk – $10B already spent is "
        "not a reason to keep going. Second, the real reason to stop "
        "was opportunity cost: the capital and the ~2,000 engineers "
        "had a higher-return use in Apple Intelligence and the AI "
        "stack. Reallocation, not refusal-to-cut-losses, is the "
        "executive lesson. Opportunity cost is the next-best "
        "alternative use of the same dollars and engineers."
    ))


def slide_48(prs):
    """Dictionary of costs – native three-card cheat sheet.

    Same formulas and labels as the source image, restyled to the deck's
    NAVY/GOLD palette with OMML rendering for each headline formula.

    User hand-edits ported 2026-05-22:
      • "Cheat sheet to refer back to" caption moved from below the
        cards (T 6.60) to ABOVE the cards (T 1.90), text shortened
        ("…for the rest of the module" dropped), font 14 → 24 pt.
      • Three cards (header rect + OMML formula box + formula
        textbox) shifted down by 0.89" to make room for the caption:
        hdr_y 1.85 → 2.74.
      • Sub-formula textbox fonts bumped from 15/12 pt → 20/18 pt
        for EMBA legibility.
    """
    def draw(slide):
        # Three-card row: Total / Average / Marginal cost.
        col_w = Inches(4.05)
        col_gap = Inches(0.15)
        col_x0 = (SLIDE_W - col_w * 3 - col_gap * 2) // 2

        # "Cheat sheet…" caption — moved ABOVE the cards (2026-05-22
        # hand-edit: was at T=6.60 below the cards, now at T=1.90
        # spanning roughly the column band). Font bumped 14 → 24 pt.
        _add_text(slide, Inches(0.44), Inches(1.90), Inches(12.78), Inches(0.40),
                  "Cheat sheet to refer back to",
                  size=24, italic=True, color=GRAY, font="Calibri",
                  align=PP_ALIGN.CENTER)

        # 2026-05-22 hand-tweaked: hdr_y 1.85 → 2.74 to clear the
        # newly-placed top caption.  OMML and sub-textbox tops flow
        # from this.
        hdr_y = Inches(2.74)
        hdr_h = Inches(0.55)
        formula_y = hdr_y + hdr_h + Inches(0.10)
        formula_h = Inches(0.95)
        subs_y = formula_y + formula_h + Inches(0.15)

        # Card 1 — Total Cost
        x = col_x0
        _add_filled_box(slide, x, hdr_y, col_w, hdr_h,
                        "Total Cost", fill=NAVY, text_color=WHITE,
                        size=22, bold=True)
        _add_math_equation(
            slide, x, formula_y, col_w, formula_h,
            _omml_text('TC = TFC + TVC'),
            size_pt=22, color=NAVY, fill=RGBColor(0xFD, 0xF6, 0xE6),
            line=NAVY,
        )
        # 2026-05-22 hand-tweak: bullet fonts 15/12 → 20/18 pt.
        _add_hierarchical_bullets(
            slide,
            left=x + Inches(0.15), top=subs_y,
            width=col_w - Inches(0.30), height=Inches(2.80),
            items=[
                ("TFC = Total Fixed Cost", 0),
                ("(ignore sunk costs)", 1),
                ("TVC = Total Variable Cost", 0),
            ],
            size=20, sub_size=18, line_spacing_pts=8,
        )

        # Card 2 — Average Cost
        x = col_x0 + col_w + col_gap
        _add_filled_box(slide, x, hdr_y, col_w, hdr_h,
                        "Average Cost", fill=NAVY, text_color=WHITE,
                        size=22, bold=True)
        _add_math_equation(
            slide, x, formula_y, col_w, formula_h,
            _omml_text('ATC = ') + _omml_frac(_omml_text('TC'), _omml_text('Q')),
            size_pt=22, color=NAVY, fill=RGBColor(0xFD, 0xF6, 0xE6),
            line=NAVY,
        )
        _add_hierarchical_bullets(
            slide,
            left=x + Inches(0.15), top=subs_y,
            width=col_w - Inches(0.30), height=Inches(2.80),
            items=[
                ("AFC = TFC / Q", 0),
                ("AVC = TVC / Q", 0),
                ("ATC = AFC + AVC", 0),
            ],
            size=20, sub_size=18, line_spacing_pts=8,
        )

        # Card 3 — Marginal Cost
        x = col_x0 + 2 * (col_w + col_gap)
        _add_filled_box(slide, x, hdr_y, col_w, hdr_h,
                        "Marginal Cost", fill=NAVY, text_color=WHITE,
                        size=22, bold=True)
        _add_math_equation(
            slide, x, formula_y, col_w, formula_h,
            _omml_text('MC = ') + _omml_frac(_omml_text('Δ') + _omml_text('TC'),
                                                _omml_text('Δ') + _omml_text('Q')),
            size_pt=22, color=NAVY, fill=RGBColor(0xFD, 0xF6, 0xE6),
            line=NAVY,
        )
        _add_hierarchical_bullets(
            slide,
            left=x + Inches(0.15), top=subs_y,
            width=col_w - Inches(0.30), height=Inches(2.80),
            items=[
                ("= ΔTVC / ΔQ  (TFC is constant)", 0),
                ("Derivative form:", 0),
                ("MC = dTC / dQ", 1),
                ("    = dTVC / dQ", 1),
            ],
            size=20, sub_size=18, line_spacing_pts=8,
        )

    s = make_diagram_slide(
        prs, page_num=48,
        section_tag=SECTION_TAG_P2,
        title="Dictionary of Costs",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "Quick reference: fixed, variable, sunk, marginal, average, "
        "opportunity. A cheat sheet you'll refer to for the rest of the "
        "module. Make sure you can give a one-sentence example of each."
    ))


def slide_49(prs):
    """Relationship to Accounting: Ross Stores annual report.

    2026-05-22 third pass: the user replaced my multi-column
    reconstruction with a single combined 10-K screenshot
    (_ross_costs_combined.png) plus four red text annotations.
    This version reproduces the user's paste verbatim:

      • Single combined cost-table image at
        L 0.82 / T 1.62 / W 11.46 / H 5.21.
      • Four red text annotations (all 20 pt bold red C00000,
        Whitney Book font — PowerPoint substitutes Calibri
        locally where Whitney Book isn't installed):
          - "≈ TVC"        at L 3.45 / T 4.55   (Cost of goods sold row)
          - "Mix FC & VC"  at L 5.18 / T 5.11   (SG&A row)
          - "Part of FC"   at L 4.93 / T 5.59   (Interest expense row)
          - "≈ TC"         at L 4.20 / T 6.19   (Total costs and expenses row)

    Coordinates copied verbatim from the user's pasted shapes.
    """
    def draw(slide):
        # ----- Single combined Ross Stores cost-table image ------
        # User pasted this as one image (replaces the previous
        # multi-column source-image reconstruction, which looked
        # distorted in 16:9).
        ross = OUT_DIR / "_ross_costs_combined.png"
        if ross.exists():
            slide.shapes.add_picture(
                str(ross),
                int(Inches(0.82)), int(Inches(1.62)),
                width=int(Inches(11.46)), height=int(Inches(5.21)),
            )

        # ----- Four red text annotations -------------------------
        # 2026-05-23 hand-nudges to fine-tune row alignment: each of the
        # four annotations was moved a few hundredths of an inch.
        RED = RGBColor(0xC0, 0x00, 0x00)
        # "≈ TVC" — aligned with the "Cost of goods sold" row.
        _add_text(slide, Inches(3.49), Inches(4.58),
                   Inches(1.16), Inches(0.44),
                   "≈ TVC",
                   size=20, bold=True, color=RED, font="Whitney Book")
        # "Mix FC & VC" — aligned with the SG&A row.
        _add_text(slide, Inches(5.26), Inches(5.12),
                   Inches(1.98), Inches(0.44),
                   "Mix FC & VC",
                   size=20, bold=True, color=RED, font="Whitney Book")
        # "Part of FC" — aligned with the Interest-expense row.
        _add_text(slide, Inches(4.92), Inches(5.66),
                   Inches(1.45), Inches(0.44),
                   "Part of FC",
                   size=20, bold=True, color=RED, font="Whitney Book")
        # "≈ TC" — aligned with the "Total costs and expenses" row.
        _add_text(slide, Inches(4.24), Inches(6.26),
                   Inches(1.16), Inches(0.44),
                   "≈ TC",
                   size=20, bold=True, color=RED, font="Whitney Book")

    s = make_diagram_slide(
        prs, page_num=49,
        section_tag=SECTION_TAG_P2,
        title="Relationship to Accounting:  Ross Stores Annual Report",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "Ross Stores' annual report, used here to ground the cost "
        "concepts in something students will recognise from a 10-K. "
        "Walk through the three items under 'costs and expenses': "
        "cost of goods sold is essentially variable (≈ TVC); "
        "selling, general and administrative is a mix of fixed and "
        "variable (sales commissions vary with revenue, the CEO's "
        "salary does not); interest expense is fixed (the cost of "
        "servicing long-term debt). Everything labelled costs and "
        "expenses together is the firm's TC. The red overlays make "
        "the mapping visible without rewriting the line items."
    ))


def slide_50(prs):
    """Marginal Cost in Action: MC of a Burn60 Workout.

    2026-05-23 third pass: the user replaced the previous build's
    image + separate "PACKAGES" pricing textbox with a single
    pre-composited screenshot in which the package pricing is
    baked into the bitmap.  The only overlay left on top of the
    image is the red ellipse circling the package block.  The
    slide now contains just three content pieces (chrome aside):
    the merged screenshot, the red circle, and the
    "Poll Everywhere break" badge in the lower-right corner.
    """
    def draw(slide):
        # ----- Merged Burn60 screenshot (pricing baked in) -------
        # 2026-05-23 hand-positioned:
        #   • Picture L 2.58 → 2.84, T 1.24 → 1.39 (image scaled
        #     down slightly; new W 6.95, H 5.71).
        #   • The standalone PACKAGES textbox is GONE — its text
        #     now lives inside the screenshot's PNG bitmap.
        burn60 = OUT_DIR / "_burn60_workout.png"
        if burn60.exists():
            slide.shapes.add_picture(
                str(burn60),
                int(Inches(2.84)), int(Inches(1.39)),
                width=int(Inches(6.95)), height=int(Inches(5.71)),
            )

        # ----- Red ellipse circling the package block -----------
        # 2026-05-23 hand-tweak: oval re-positioned to track the
        # smaller image (L 4.70 → 4.74, T 5.24 → 5.37).
        RED = RGBColor(0xFF, 0x00, 0x00)
        oval = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            int(Inches(4.74)), int(Inches(5.37)),
            int(Inches(2.98)), int(Inches(1.51)),
        )
        oval.fill.background()
        oval.line.color.rgb = RED
        oval.line.width = Pt(3.0)

        # ----- "Poll Everywhere break" badge (bottom-right) -----
        _add_discussion_break(
            slide,
            top=Inches(6.34),
            width=Inches(4.8),
            text="Poll Everywhere break",
        )

    s = make_diagram_slide(
        prs, page_num=50,
        section_tag=SECTION_TAG_P2,
        title="Marginal Cost in Action:  What Is the MC of a Burn60 Workout?",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "Real-world setup for the marginal-cost concept. Burn60 (a "
        "boutique fitness studio) offers package pricing: 10 classes "
        "for $10 each ($100 total) or 20 classes for $8 each "
        "($160 total). The question to put to the room: starting "
        "from the 10-pack, what's the marginal cost of the 11th "
        "class? Hint: it's not $8. Going from the 10-pack to the "
        "20-pack costs an additional $60 for ten more classes — so "
        "the marginal cost of each of those next ten is $6, not the "
        "$8 sticker price. Average cost (per-unit price you see "
        "advertised) and marginal cost (what one more actually "
        "costs) diverge whenever there is volume-tier pricing. We "
        "poll the room on the next slide before walking through it."
    ))


def slide_51(prs):
    """Poll: MC of each additional Burn60 session beyond the 10th?

    Source slide is a full-bleed PollEv screenshot showing the Burn60
    marginal-cost question with four candidate answers ($10, $8, $18,
    $6) plus "None of the above".  2026-05-23 retitled and re-chromed
    from the prior ChatGPT framing.
    """
    def draw(slide):
        _add_source_image(slide, 52, "rId4",
                          left=Inches(3.2), top=Inches(1.85),
                          height=Inches(5.1))
        _add_text(slide, MARGIN, Inches(7.0), RULE_W, Inches(0.3),
                  "Respond at PollEv.com/nvoigtlaender",
                  size=14, italic=True, color=GRAY,
                  align=PP_ALIGN.CENTER, font="Calibri")
        # 2026-05-23: replaced the top-right POLL pill (yellow
        # rectangle + label) with the deck-standard gold
        # "Poll Everywhere break" badge in the lower-right corner,
        # matching slide 50.
        _add_discussion_break(
            slide,
            top=Inches(6.34),
            width=Inches(4.8),
            text="Poll Everywhere break",
        )

    s = make_diagram_slide(
        prs, page_num=51,
        section_tag=SECTION_TAG_P2,
        title="What Is the MC of Each Additional Session Beyond the 10th?",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "Quick PollEv.  Compute the marginal cost of each additional "
        "Burn60 session when buying the 20-pack instead of the 10-pack.  "
        "The common trap: students see the $8 per-class rate on the "
        "20-pack and answer $8.  But going from the 10-pack ($100 "
        "total) to the 20-pack ($160 total) costs $60 more for 10 "
        "more sessions — so the MC of each additional session is $6, "
        "not $8.  The first 10 sessions get re-priced from $10 to $8 "
        "as well; that re-pricing is what pulls the marginal cost "
        "below the headline 20-pack rate.  Give them 30 seconds, then "
        "reveal the solution on the next slide."
    ))


def slide_52(prs):
    """Solution: MC = $6 / session.

    2026-05-23 rebuild: dropped the ChatGPT Plus-vs-Team solution at
    the user's direction; mirrors the original deck's slide 58 (Burn60
    package pricing).  Layout primitives kept identical to the prior
    ChatGPT solution slide — two side-by-side cost-line boxes, a
    centered MC-formula band, intuition line, and a takeaway bar.
    """
    def draw(slide):
        # Two side-by-side cost lines
        col_w = Inches(5.6)
        col_gap = Inches(0.5)
        x0 = (SLIDE_W - col_w * 2 - col_gap) // 2

        _add_filled_box(
            slide, x0, Inches(1.95), col_w, Inches(0.55),
            "10-class pack  ($10 / class)", fill=NAVY, text_color=WHITE,
            size=20, bold=True,
        )
        _add_outlined_box(
            slide, x0, Inches(2.5), col_w, Inches(0.9),
            "TC  =  10 × $10  =  $100",
            fill=WHITE, line=NAVY, text_color=NAVY,
            size=22, bold=True, line_w=1.5,
        )

        _add_filled_box(
            slide, x0 + col_w + col_gap, Inches(1.95), col_w, Inches(0.55),
            "20-class pack  ($8 / class)", fill=GOLD, text_color=NAVY,
            size=20, bold=True,
        )
        _add_outlined_box(
            slide, x0 + col_w + col_gap, Inches(2.5), col_w, Inches(0.9),
            "TC  =  20 × $8  =  $160",
            fill=WHITE, line=GOLD, text_color=NAVY,
            size=22, bold=True, line_w=1.5,
        )

        # MC calculation row
        # 2026-05-23 hand-tweak: header font 22 → 28 pt.
        _add_text(slide, MARGIN, Inches(3.7), RULE_W, Inches(0.5),
                  "Marginal cost of each additional session beyond the 10th:",
                  size=28, italic=True, color=GRAY, font="Calibri",
                  align=PP_ALIGN.CENTER)
        _add_filled_box(
            slide, (SLIDE_W - Inches(8.5)) // 2, Inches(4.25),
            Inches(8.5), Inches(1.0),
            "MC  =  ($160 − $100) / 10  =  $6 / session",
            fill=NAVY, text_color=WHITE, size=28, bold=True,
        )

        # Intuition
        _add_text(slide, MARGIN, Inches(5.5), RULE_W, Inches(0.5),
                  "Bigger discount than the sticker:  the first 10 sessions also drop from $10 → $8",
                  size=18, italic=True, color=NAVY, font="Calibri",
                  align=PP_ALIGN.CENTER)

        # 2026-05-23 hand-tweaks: takeaway bar moved up (T 6.5 → 6.22)
        # and rephrased — "volume pricing" → "volume discounts".
        _add_takeaway_bar(
            slide,
            "MC  <  AC :  volume discounts make the marginal session cheaper than the sticker",
            top=Inches(6.22), fill=GOLD, text_color=NAVY,
            width=Inches(11.6),
        )

    s = make_diagram_slide(
        prs, page_num=52,
        section_tag=SECTION_TAG_P2,
        title="Solution:  MC = $6 / Session",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "Reveal: MC = $6 per additional session — not the $8 sticker "
        "rate on the 20-pack. Total cost goes from $100 (10 classes at "
        "$10) to $160 (20 classes at $8), so the next ten cost $60 "
        "incremental → $6 per session. The intuition: the 20-pack also "
        "re-prices the first ten sessions from $10 down to $8, and "
        "that $2 × 10 = $20 of savings on the inframarginal sessions "
        "is what pulls the marginal rate below the $8 sticker. "
        "Pedagogical point: volume pricing makes the marginal unit "
        "cheaper than the headline per-unit rate — the mirror image "
        "of the SaaS-tier example where the marginal user is more "
        "expensive than the headline rate."
    ))


def slide_53(prs):
    """Marginal cost in finance: the true rate on a bigger loan.

    2026-05-23 full redesign per user mockup:
      • Two Option cards side-by-side — navy header band ("Option 1" /
        "Option 2") over a white rounded body with a navy border and
        soft drop shadow.
      • Each card body: a navy HOME-shape bank icon, a large loan
        title ("$100k loan" / "$110k loan", 28 pt bold navy) and an
        italic subtitle ("at 5% annual interest" / "at 6% annual
        interest", 18 pt italic navy).
      • Three labelled rows per card below a divider line, each a
        SEPARATE text box (so the user can fly them in independently
        when animating).  Each row carries a small navy-outlined
        icon (%, ≡, $) on the left, the label, and the right-aligned
        value via a right-aligned tab stop.
      • A red dashed arrow between the two cards labelled
        "Extra $10k".
      • A cream/peach rounded calculation box under the cards: the
        question ("What is the marginal interest rate on the
        additional $10k?", red bold) on line 1, the worked
        calculation ("($6,600 − $5,000) / $10,000 = 16%", navy
        bold) on line 2.
      • The existing yellow takeaway bar at the bottom — now made
        rounded + drop-shadowed and narrowed from W 11.0 to 10.0.
    """
    def draw(slide):
        from pptx.oxml.ns import qn

        # ----- Card geometry -----
        # 2026-05-23 second-pass user tweaks: cards narrowed and the
        # gap between them widened so the "Extra $10k" arrow has
        # significantly more room (W 5.50 → 4.80, gap 0.55 → 1.50).
        card_w = Inches(4.80)
        card_gap = Inches(1.50)
        card_x0 = (SLIDE_W - 2 * card_w - card_gap) // 2
        header_h = Inches(0.55)
        body_h = Inches(3.30)
        card_t = Inches(1.55)
        body_t = card_t + header_h

        RED = RGBColor(0xC0, 0x00, 0x00)
        CREAM = RGBColor(0xFD, 0xEC, 0xDB)
        DIVIDER = RGBColor(0xC8, 0xCE, 0xD6)

        def build_card(x, option_label, loan_amount, interest_pct,
                       annual_interest, loan_total, interest_cost):
            # ---- Rounded navy header ("Option 1" / "Option 2") ----
            # User asked for rounded corners + drop shadow on the
            # Option headers (was a flat filled rectangle).
            _add_rounded_filled_box(
                slide, x, card_t, card_w, header_h,
                label=option_label,
                fill=NAVY, text_color=WHITE,
                size=22, bold=True,
                corner_pct=0.20, shadow=True,
            )
            # ---- Rounded white body card (navy border + shadow) ----
            body = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                int(x), int(body_t), int(card_w), int(body_h),
            )
            try: body.adjustments[0] = 0.04
            except Exception: pass
            body.fill.solid(); body.fill.fore_color.rgb = WHITE
            body.line.color.rgb = NAVY; body.line.width = Pt(0.75)
            body.shadow.inherit = False
            _add_drop_shadow(body)

            # ---- House / bank pictogram (top-left of body) ----
            # MSO_SHAPE.HOME isn't available in this python-pptx, so
            # we compose a simple house silhouette from two primitives:
            # a navy filled rectangle (building body) sitting under a
            # navy filled isoceles triangle (pitched roof).
            icon_x = x + Inches(0.25)
            icon_y = body_t + Inches(0.30)
            roof = slide.shapes.add_shape(
                MSO_SHAPE.ISOSCELES_TRIANGLE,
                int(icon_x), int(icon_y),
                int(Inches(0.70)), int(Inches(0.28)),
            )
            roof.fill.solid(); roof.fill.fore_color.rgb = NAVY
            roof.line.fill.background()
            body_rect = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                int(icon_x + Inches(0.05)),
                int(icon_y + Inches(0.26)),
                int(Inches(0.60)), int(Inches(0.34)),
            )
            body_rect.fill.solid(); body_rect.fill.fore_color.rgb = NAVY
            body_rect.line.fill.background()

            # ---- Large loan title + italic subtitle ----
            # 2026-05-23 second-pass: subtitle 18 → 22 pt.
            title_x = x + Inches(1.05)
            _add_text(slide, title_x, body_t + Inches(0.20),
                      card_w - Inches(1.20), Inches(0.50),
                      loan_amount, size=28, bold=True, color=NAVY,
                      font="Calibri")
            _add_text(slide, title_x, body_t + Inches(0.72),
                      card_w - Inches(1.20), Inches(0.45),
                      f"at {interest_pct} annual interest",
                      size=22, italic=True, color=NAVY, font="Calibri")

            # ---- Divider rule under title ----
            div = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                int(x + Inches(0.25)),
                int(body_t + Inches(1.35)),
                int(card_w - Inches(0.50)),
                int(Inches(0.012)),
            )
            div.fill.solid(); div.fill.fore_color.rgb = DIVIDER
            div.line.fill.background()

            # ---- Three labelled rows ----
            # 2026-05-23 second-pass: row text 15 → 20 pt; row height
            # 0.40 → 0.50 to accommodate; icons scaled 0.32 → 0.40.
            row_h = Inches(0.50)
            row_gap = Inches(0.10)
            row_y0 = body_t + Inches(1.50)
            icon_size = Inches(0.40)
            label_x = x + Inches(0.80)
            label_w = card_w - Inches(1.05)

            def build_row(idx, icon_char, label, value):
                row_y = row_y0 + idx * (row_h + row_gap)
                # --- Icon: small navy-outlined oval with letter ---
                ic = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    int(x + Inches(0.28)),
                    int(row_y + (row_h - icon_size) // 2),
                    int(icon_size), int(icon_size),
                )
                ic.fill.solid(); ic.fill.fore_color.rgb = WHITE
                ic.line.color.rgb = NAVY; ic.line.width = Pt(1.25)
                itf = ic.text_frame
                itf.margin_left = 0; itf.margin_right = 0
                itf.margin_top = 0; itf.margin_bottom = 0
                itf.vertical_anchor = MSO_ANCHOR.MIDDLE
                ip = itf.paragraphs[0]; ip.alignment = PP_ALIGN.CENTER
                ir = ip.add_run(); ir.text = icon_char
                ir.font.name = "Calibri"; ir.font.size = Pt(16)
                ir.font.bold = True; ir.font.color.rgb = NAVY

                # --- Row text box: label left + value right (tab) ---
                box = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    int(label_x), int(row_y),
                    int(label_w), int(row_h),
                )
                box.fill.background(); box.line.fill.background()
                btf = box.text_frame
                btf.margin_left = 0; btf.margin_right = Inches(0.05)
                btf.margin_top = 0; btf.margin_bottom = 0
                btf.vertical_anchor = MSO_ANCHOR.MIDDLE
                p = btf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
                # Right-aligned tab stop at the right edge.
                pPr = p._p.get_or_add_pPr()
                tabLst = ET.SubElement(pPr, qn('a:tabLst'))
                tab = ET.SubElement(tabLst, qn('a:tab'))
                tab.set('pos', str(int(label_w)))
                tab.set('algn', 'r')
                # Label run (regular)
                r1 = p.add_run(); r1.text = label
                r1.font.name = "Calibri"; r1.font.size = Pt(20)
                r1.font.color.rgb = NAVY
                # Value run (bold) preceded by tab
                r2 = p.add_run(); r2.text = "\t" + value
                r2.font.name = "Calibri"; r2.font.size = Pt(20)
                r2.font.bold = True; r2.font.color.rgb = NAVY

            build_row(0, "%", "Annual interest", annual_interest)
            build_row(1, "≡", "Total loan amount", loan_total)
            build_row(2, "$", "Annual interest cost", interest_cost)

        # Build both Option cards
        build_card(
            card_x0, "Option 1", "$100k loan", "5%",
            "$5,000", "$100,000", "$5,000",
        )
        build_card(
            card_x0 + card_w + card_gap, "Option 2", "$110k loan", "6%",
            "$6,600", "$110,000", "$6,600",
        )

        # ----- "Extra $10k" dashed arrow between cards -----
        # 2026-05-23 second-pass: label significantly larger
        # (13 → 24 pt bold) and arrow thickened (2.0 → 4.0 pt) so the
        # callout reads from the back row.
        arrow_y = body_t + Inches(0.75)
        arrow_x_start = card_x0 + card_w + Inches(0.10)
        arrow_x_end = card_x0 + card_w + card_gap - Inches(0.06)
        # Label above arrow
        _add_text(slide,
                   arrow_x_start - Inches(0.10),
                   arrow_y - Inches(0.65),
                   card_gap, Inches(0.50),
                   "Extra $10k",
                   size=24, bold=True, color=RED, font="Calibri",
                   align=PP_ALIGN.CENTER)
        # Red dashed arrow — thicker, more prominent
        _add_arrow(slide,
                   (arrow_x_start, arrow_y),
                   (arrow_x_end, arrow_y),
                   color=RED, weight_pt=4.0, dash='dash')

        # ----- Cream calculation box (question + computation) -----
        # 2026-05-23 second-pass: question font 18 → 26 pt; calc box
        # height bumped 0.95 → 1.20 in to fit the larger text.
        calc_w = Inches(10.80)
        calc_left = (SLIDE_W - calc_w) // 2
        calc_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            int(calc_left), int(Inches(5.45)),
            int(calc_w), int(Inches(1.00)),
        )
        try: calc_box.adjustments[0] = 0.18
        except Exception: pass
        calc_box.fill.solid(); calc_box.fill.fore_color.rgb = CREAM
        calc_box.line.color.rgb = RED; calc_box.line.width = Pt(0.75)
        calc_box.shadow.inherit = False
        _add_drop_shadow(calc_box)
        ctf = calc_box.text_frame
        ctf.margin_left = Inches(0.30); ctf.margin_right = Inches(0.20)
        ctf.margin_top = Inches(0.10); ctf.margin_bottom = Inches(0.05)
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        # Line 1: question (red bold) — significantly larger
        p0 = ctf.paragraphs[0]; p0.alignment = PP_ALIGN.CENTER
        r0 = p0.add_run()
        r0.text = "What is the marginal interest rate on the additional $10k?"
        r0.font.name = "Calibri"; r0.font.size = Pt(26)
        r0.font.bold = True; r0.font.color.rgb = RED
        # Line 2: calculation (navy bold)
        p1 = ctf.add_paragraph(); p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = "($6,600 − $5,000) / $10,000 = 16%"
        r1.font.name = "Calibri"; r1.font.size = Pt(24)
        r1.font.bold = True; r1.font.color.rgb = NAVY

        # ----- Yellow takeaway bar (rounded + shadow, narrower) -----
        _add_rounded_filled_box(
            slide,
            (SLIDE_W - Inches(10.0)) // 2, Inches(6.55),
            Inches(10.0), Inches(0.50),
            label="Marginal cost (16%) is much higher than average rate (6%)",
            fill=GOLD, text_color=NAVY,
            size=18, bold=True,
            corner_pct=0.30, shadow=True,
        )

    s = make_diagram_slide(
        prs, page_num=53,
        section_tag=SECTION_TAG_P2,
        title="Marginal Cost in Finance:  The True Rate on a Bigger Loan",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "Same concept applied to finance.  Option 1: a $100k loan at "
        "5% costs $5,000 in interest per year.  Option 2: a $110k loan "
        "at 6% costs $6,600 in interest per year.  The marginal "
        "interest rate on the extra $10k is the EXTRA interest divided "
        "by the EXTRA principal — ($6,600 − $5,000) / $10,000 = 16%.  "
        "The headline rate of 6% is the AVERAGE; the marginal rate "
        "you actually pay for the extra dollars is 16%.  This is the "
        "common executive trap when comparing loan terms — the rate "
        "advertised on the larger loan understates what the marginal "
        "dollar truly costs."
    ))


def slide_54(prs):
    """Rivian's Georgia plant – weekly Total Cost function.

    2026-05-23 redesign to mirror original deck slide 60:
      • Top intro text bumped 22 → 28 pt (user hand-edit).
      • TC formula switched to the original deck's
        TC = 10,000,000 + 30,000·Q + 40·Q² (was 800k + 200Q²).
      • Native chart redrawn as XY scatter: an "Observed" series
        of points scattered around the true curve (~6 % gaussian
        noise, fixed seed) plus a smooth "Fitted" line drawing
        the true TC function.  Replaces the prior perfect-fit
        line chart.
      • "Fixed Cost" annotation with red arrow pointing to the
        y-intercept ($10M at Q=0).
      • Y axis re-scaled to $M (max 90) so labels stay clean
        across the larger value range.

    Formula constants are LOCAL to this slide so slides 55-57
    (which still use the 800k + 200Q² version via COST_TFC /
    COST_VAR_COEF / COST_Q_VALS) keep building.  Propagate to
    those slides only if/when the user asks.
    """
    def draw(slide):
        from pptx.chart.data import XyChartData
        from pptx.enum.chart import XL_MARKER_STYLE

        # ----- Condensed top intro (one line) -------------------
        # 2026-05-25 hand-tweaks: text bumped 22 → 24 pt italic;
        # box moved up T 1.55 → 1.46, height tightened 0.45 → 0.40.
        _add_text(slide, MARGIN, Inches(1.46), RULE_W, Inches(0.40),
                  "Collect data on total cost at different output levels (Q)  →  "
                  "Estimate the cost function",
                  size=24, italic=True, color=NAVY, font="Calibri",
                  align=PP_ALIGN.CENTER)

        # ----- Cost function (OMML) ------------------------------
        # 2026-05-24 user updates:
        #   • Prefix the equation with "Estimated cost curve:".
        #   • Move the box up (T 2.10 → 1.95) so the chart has more
        #     vertical room below.
        # 2026-05-25: equation box gets rounded corners + soft drop
        # shadow to match the deck's hero-formula treatment.
        eq_xml = (
            _omml_text('Estimated cost curve:   ') +
            _omml_text('TC') +
            _omml_text(' = ') +
            _omml_text('10,000,000') +
            _omml_text(' + ') +
            _omml_text('30,000') +
            _omml_text(' · ') +
            _omml_run('Q') +
            _omml_text(' + ') +
            _omml_text('40') +
            _omml_text(' · ') +
            _omml_sup(_omml_run('Q'), _omml_text('2'))
        )
        eq_box = _add_math_equation(
            slide, (SLIDE_W - Inches(11.0)) // 2, Inches(1.95),
            Inches(11.0), Inches(0.60),
            eq_xml, size_pt=22, color=WHITE, fill=NAVY,
        )
        # Swap the textbox's rectangle geometry for a rounded one
        # and attach a soft drop shadow.
        # IMPORTANT: in OOXML, <a:prstGeom> must come immediately
        # after <a:xfrm> and BEFORE any fill / effect element.
        # If we append it with SubElement, it lands AFTER the
        # solidFill that _add_math_equation set, and PowerPoint
        # silently refuses to render the shape (it "disappears").
        eq_spPr = eq_box._element.find(qn('p:spPr'))
        if eq_spPr is not None:
            for old in eq_spPr.findall(qn('a:prstGeom')):
                eq_spPr.remove(old)
            prstGeom = ET.Element(qn('a:prstGeom'))
            prstGeom.set('prst', 'roundRect')
            avLst = ET.SubElement(prstGeom, qn('a:avLst'))
            gd = ET.SubElement(avLst, qn('a:gd'))
            gd.set('name', 'adj'); gd.set('fmla', 'val 25000')
            # Position prstGeom right after a:xfrm (schema-correct).
            xfrm = eq_spPr.find(qn('a:xfrm'))
            if xfrm is not None:
                xfrm.addnext(prstGeom)
            else:
                eq_spPr.insert(0, prstGeom)
        _add_drop_shadow(eq_box)

        # ----- Local TC formula ----------------------------------
        TFC_LOCAL = 10_000_000
        LIN_LOCAL = 30_000
        QUAD_LOCAL = 40
        tc_of = lambda q: TFC_LOCAL + LIN_LOCAL * q + QUAD_LOCAL * q * q

        # ----- Observed scatter ---------------------------------
        # 2026-05-24 (second pass) user request: the previous
        # hand-tuned offsets alternated signs and read as an
        # artificial "up-down-up-down" pattern.  Switch to truly
        # random uniform deviations in [-0.20, +0.20] (±20% on
        # TC), rejection-sampled so no offset is smaller than 6 %
        # (otherwise some dots would sit visually on the line).
        # Seed = 3 chosen because it gives a balanced spread of
        # signs (no run longer than 3) and magnitudes 0.10–0.20.
        import random
        random.seed(3)
        Q_obs = list(range(0, 1001, 100))
        obs_offsets = []
        for _ in Q_obs:
            while True:
                o = random.uniform(-0.20, 0.20)
                if abs(o) >= 0.06:
                    break
            obs_offsets.append(o)
        TC_obs_M = [tc_of(q) * (1.0 + off) / 1_000_000
                    for q, off in zip(Q_obs, obs_offsets)]
        # Fitted curve: dense points (every 25 units) so the line
        # connecting them reads as smooth.
        Q_curve = list(range(0, 1001, 25))
        TC_curve_M = [tc_of(q) / 1_000_000 for q in Q_curve]

        # ----- Chart data ----------------------------------------
        # 2026-05-24: legend added, series renamed to "TC" (line)
        # and "TC (data)" (markers).  Line series goes FIRST so
        # the data dots are drawn ON TOP and visible.
        chart_data = XyChartData()
        s_fit = chart_data.add_series('TC')
        for q, tc in zip(Q_curve, TC_curve_M):
            s_fit.add_data_point(q, tc)
        s_obs = chart_data.add_series('TC (data)')
        for q, tc in zip(Q_obs, TC_obs_M):
            s_obs.add_data_point(q, tc)

        # Chart geometry: taller now that the equation is higher.
        # T 2.85 → 2.65;  H 3.60 → 3.80.
        chart_x, chart_y = Inches(2.50), Inches(2.65)
        chart_w, chart_h = Inches(8.30), Inches(3.80)
        gf = slide.shapes.add_chart(
            XL_CHART_TYPE.XY_SCATTER_LINES,
            chart_x, chart_y, chart_w, chart_h, chart_data,
        )
        chart = gf.chart

        # ----- Axes (larger tick labels + titles) ----------------
        xax = chart.category_axis
        xax.maximum_scale = 1000
        xax.minimum_scale = 0
        xax.major_unit = 200
        xax.has_title = True
        xax.axis_title.text_frame.text = "Q   (vehicles per week)"
        for r in xax.axis_title.text_frame.paragraphs[0].runs:
            r.font.size = Pt(16); r.font.italic = True; r.font.bold = True
            r.font.color.rgb = NAVY; r.font.name = "Calibri"
        xax.tick_labels.font.size = Pt(14)
        xax.tick_labels.font.color.rgb = NAVY
        xax.tick_labels.font.name = "Calibri"

        yax = chart.value_axis
        yax.maximum_scale = 110
        yax.minimum_scale = 0
        yax.major_unit = 10
        yax.has_title = True
        yax.axis_title.text_frame.text = "TC   ($M)"
        for r in yax.axis_title.text_frame.paragraphs[0].runs:
            r.font.size = Pt(16); r.font.italic = True; r.font.bold = True
            r.font.color.rgb = NAVY; r.font.name = "Calibri"
        yax.tick_labels.font.size = Pt(14)
        yax.tick_labels.font.color.rgb = NAVY
        yax.tick_labels.font.name = "Calibri"

        # ----- Dashed light-gray major gridlines on the Y axis ---
        # Use the shared helper for consistency with the other
        # Cobb-Douglas-style charts in the deck.
        _add_dashed_gridlines(yax._element)

        # ----- Per-series styling --------------------------------
        # Series 0 (TC = fitted line): smooth line, no markers.
        s_line = chart.plots[0].series[0]
        s_line.marker.style = XL_MARKER_STYLE.NONE
        s_line.format.line.color.rgb = NAVY
        s_line.format.line.width = Pt(2.75)
        s_line_el = s_line._element
        for old in s_line_el.findall(qn('c:smooth')):
            s_line_el.remove(old)
        smooth = ET.SubElement(s_line_el, qn('c:smooth'))
        smooth.set('val', '1')

        # Series 1 (TC (data) = observed markers): circles only,
        # hide connecting line.
        s_dots = chart.plots[0].series[1]
        s_dots.marker.style = XL_MARKER_STYLE.CIRCLE
        s_dots.marker.size = 9
        try:
            s_dots.marker.format.fill.solid()
            s_dots.marker.format.fill.fore_color.rgb = NAVY
            s_dots.marker.format.line.color.rgb = NAVY
        except Exception:
            pass
        s_dots_el = s_dots._element
        spPr_dots = s_dots_el.find(qn('c:spPr'))
        if spPr_dots is None:
            spPr_dots = ET.SubElement(s_dots_el, qn('c:spPr'))
        for child in list(spPr_dots):
            if child.tag == qn('a:ln'):
                spPr_dots.remove(child)
        ln_dots = ET.SubElement(spPr_dots, qn('a:ln'))
        ET.SubElement(ln_dots, qn('a:noFill'))

        # ----- Legend: top-left inside the plot area -----------
        # 2026-05-25 user request: bump the legend font size and
        # stack the two entries vertically (one below the other).
        # A narrow-tall manual-layout box forces single-column
        # rendering even though the deck's legend uses default
        # auto-flow.  2026-05-26 follow-up: add a white background
        # + thin navy border, and nudge the legend a bit toward
        # the center of the chart (x 0.12 → 0.20).
        chart.has_legend = True
        leg = chart.legend
        leg.include_in_layout = False
        leg.font.size = Pt(18)
        leg.font.color.rgb = NAVY
        leg.font.name = "Calibri"
        leg_el = leg._element
        for old in leg_el.findall(qn('c:layout')):
            leg_el.remove(old)
        layout = ET.SubElement(leg_el, qn('c:layout'))
        manual = ET.SubElement(layout, qn('c:manualLayout'))
        for name, val in (('xMode', 'edge'), ('yMode', 'edge')):
            el = ET.SubElement(manual, qn('c:' + name))
            el.set('val', val)
        for name, val in (('x', '0.20'), ('y', '0.05'),
                          ('w', '0.18'), ('h', '0.25')):
            el = ET.SubElement(manual, qn('c:' + name))
            el.set('val', val)
        # White fill + thin navy border around the legend so it
        # reads as a self-contained badge rather than floating
        # text overlapping the gridlines.
        for old in leg_el.findall(qn('c:spPr')):
            leg_el.remove(old)
        leg_spPr = ET.SubElement(leg_el, qn('c:spPr'))
        leg_fill = ET.SubElement(leg_spPr, qn('a:solidFill'))
        leg_fill_rgb = ET.SubElement(leg_fill, qn('a:srgbClr'))
        leg_fill_rgb.set('val', 'FFFFFF')
        leg_ln = ET.SubElement(leg_spPr, qn('a:ln'))
        leg_ln.set('w', '9525')  # 0.75 pt
        leg_ln_fill = ET.SubElement(leg_ln, qn('a:solidFill'))
        leg_ln_rgb = ET.SubElement(leg_ln_fill, qn('a:srgbClr'))
        leg_ln_rgb.set('val', '0B2B4E')  # navy

        # ----- "Fixed Cost" annotation pointing at TC = 10 -------
        # 2026-05-25 hand-positioned by user:
        #   • Text box: L 0.55 → 0.34
        #   • Arrow Y:  5.45 → 5.335 (lands closer to the y = 10
        #     gridline now that we can see the rendered chart).
        #   • Arrow X:  start 2.60 → 2.41, end 3.16 → 2.97.
        fc_y = Inches(5.335)
        _add_text(slide,
                   Inches(0.34), fc_y - Inches(0.169),
                   Inches(2.00), Inches(0.42),
                   "Fixed Cost",
                   size=18, bold=True, color=RGBColor(0xC0, 0x00, 0x00),
                   font="Calibri", align=PP_ALIGN.RIGHT)
        _add_arrow(slide,
                    (Inches(2.41), fc_y),
                    (Inches(2.97), fc_y),
                    color=RGBColor(0xC0, 0x00, 0x00), weight_pt=2.25)

        # 2026-05-25 user hand-edit: moved up T 6.55 → 6.45 and
        # promoted to rounded corners + soft drop shadow (matches
        # the equation box above for deck-wide consistency).
        _add_takeaway_bar(
            slide,
            "Quadratic term:  cost rises faster than output, i.e. marginal cost is increasing",
            top=Inches(6.45), fill=GOLD, text_color=NAVY,
            width=Inches(10.5),
            rounded=True, shadow=True,
        )

    s = make_diagram_slide(
        prs, page_num=54,
        section_tag=SECTION_TAG_P2,
        title="Rivian's Georgia Plant —  Weekly Cost",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "Now back to Rivian.  Weekly total cost as a function of "
        "output Q for the Georgia plant.  Fixed cost of $10M, a "
        "linear variable component of $30,000 per vehicle, plus a "
        "small quadratic term ($40·Q²) that captures the rising "
        "marginal cost as production scales up.  The scattered dots "
        "are observed weekly cost data at different Q levels; the "
        "smooth navy curve is the fitted TC function — the same "
        "TC = 10,000,000 + 30,000·Q + 40·Q² used in the original "
        "deck.  Note where the curve crosses Q = 0: that y-intercept "
        "is the fixed cost (~$10M)."
    ))


def slide_55(prs):
    """Rivian's Georgia plant – Cost Components (TC = TFC + TVC).

    2026-05-23 second pass: chart fully rebuilt to mirror slide 54's
    geometry and formula.  Previously this slide used a categorical
    line chart driven by the OLD formula TC = 800k + 200·Q² with
    Q in 10..110; that drifted out of sync once slide 54 moved to
    the original-deck formula and a 0..1000 Q range.  Now:

      • XY-scatter chart (matches slide 54), Q axis 0..1000.
      • TFC = $10M constant horizontal line; TVC = 30,000·Q + 40·Q²
        rising convex curve; TC = TFC + TVC.  Y axis in $M, max
        110 (same scale as slide 54).
      • Three smooth navy/gold/red lines, no markers.
      • Narrower frame (W 12.30 → 8.30) and taller (H 4.55 → 4.90)
        with chart_title hidden so the plot area fills the frame.
      • Larger fonts everywhere: 16 pt axis titles, 14 pt tick
        labels, 18 pt legend.
      • Legend overlay inside the plot at top-left, three entries
        stacked vertically (narrow-tall manual layout), white fill
        + navy border (matching slide 54's legend treatment).
    """
    def draw(slide):
        from pptx.chart.data import XyChartData
        from pptx.enum.chart import XL_MARKER_STYLE

        # ----- TC formula (matches slide 54) ---------------------
        TFC_LOCAL = 10_000_000
        LIN_LOCAL = 30_000
        QUAD_LOCAL = 40
        tc_of  = lambda q: TFC_LOCAL + LIN_LOCAL * q + QUAD_LOCAL * q * q
        tvc_of = lambda q:             LIN_LOCAL * q + QUAD_LOCAL * q * q

        # Dense Q grid so the three curves look smooth.
        Q_curve = list(range(0, 1001, 25))
        TC_M  = [tc_of(q)  / 1_000_000 for q in Q_curve]
        TFC_M = [TFC_LOCAL / 1_000_000 for _ in Q_curve]   # constant 10
        TVC_M = [tvc_of(q) / 1_000_000 for q in Q_curve]

        chart_data = XyChartData()
        for name, vals in (('TC', TC_M), ('TFC', TFC_M), ('TVC', TVC_M)):
            s = chart_data.add_series(name)
            for q, v in zip(Q_curve, vals):
                s.add_data_point(q, v)

        # ----- Chart geometry: narrow, tall, no chart-title gap --
        # 2026-05-23 third pass: user hand-positioned the chart
        # frame to T 1.58 / W 8.27 / H 4.67 (tweaks of a few
        # hundredths of an inch each).
        chart_x, chart_y = Inches(2.50), Inches(1.58)
        chart_w, chart_h = Inches(8.27), Inches(4.67)
        # White backing rectangle for the soft drop-shadow effect.
        _add_graphicframe_shadow(slide, chart_x, chart_y, chart_w, chart_h)
        gf = slide.shapes.add_chart(
            XL_CHART_TYPE.XY_SCATTER_LINES,
            chart_x, chart_y, chart_w, chart_h, chart_data,
        )
        chart = gf.chart
        chart.has_title = False

        # ----- Axes (larger tick labels + titles) ----------------
        # 2026-05-23 second pass: axis titles 16 → 18 pt and tick
        # labels 14 → 16 pt per user request.
        xax = chart.category_axis
        xax.maximum_scale = 1000
        xax.minimum_scale = 0
        xax.major_unit = 200
        xax.has_title = True
        xax.axis_title.text_frame.text = "Q   (vehicles per week)"
        for r in xax.axis_title.text_frame.paragraphs[0].runs:
            r.font.size = Pt(18); r.font.italic = True; r.font.bold = True
            r.font.color.rgb = NAVY; r.font.name = "Calibri"
        xax.tick_labels.font.size = Pt(16)
        xax.tick_labels.font.color.rgb = NAVY
        xax.tick_labels.font.name = "Calibri"

        yax = chart.value_axis
        yax.maximum_scale = 90
        yax.minimum_scale = 0
        yax.major_unit = 10
        yax.has_title = True
        yax.axis_title.text_frame.text = "Cost   ($M)"
        for r in yax.axis_title.text_frame.paragraphs[0].runs:
            r.font.size = Pt(18); r.font.italic = True; r.font.bold = True
            r.font.color.rgb = NAVY; r.font.name = "Calibri"
        yax.tick_labels.font.size = Pt(16)
        yax.tick_labels.font.color.rgb = NAVY
        yax.tick_labels.font.name = "Calibri"

        # Dashed light-grey major gridlines on the Y axis.
        _add_dashed_gridlines(yax._element)

        # ----- Per-series styling: smooth lines, no markers -----
        colors = [NAVY, GOLD, RGBColor(0xC0, 0x50, 0x4D)]
        for i, color in enumerate(colors):
            ser = chart.plots[0].series[i]
            ser.marker.style = XL_MARKER_STYLE.NONE
            ser.format.line.color.rgb = color
            ser.format.line.width = Pt(2.75)
            ser_el = ser._element
            for old in ser_el.findall(qn('c:smooth')):
                ser_el.remove(old)
            smooth = ET.SubElement(ser_el, qn('c:smooth'))
            smooth.set('val', '1')

        # ----- Legend: top-left inside the plot area -----------
        # Vertical stacking (three entries on three lines), white
        # fill + thin navy border, 18 pt navy text.
        chart.has_legend = True
        leg = chart.legend
        leg.include_in_layout = False
        leg.font.size = Pt(18)
        leg.font.color.rgb = NAVY
        leg.font.name = "Calibri"
        leg_el = leg._element
        for old in leg_el.findall(qn('c:layout')):
            leg_el.remove(old)
        layout = ET.SubElement(leg_el, qn('c:layout'))
        manual = ET.SubElement(layout, qn('c:manualLayout'))
        for name, val in (('xMode', 'edge'), ('yMode', 'edge')):
            el = ET.SubElement(manual, qn('c:' + name))
            el.set('val', val)
        # 2026-05-23 third pass: user nudged the legend up and
        # slightly right (x 0.168 → 0.193, y 0.226 → 0.149); box
        # dimensions kept at 0.15 × 0.221 — narrow + tight enough
        # for 3 entries to stack vertically, close together.
        for name, val in (('x', '0.193'), ('y', '0.149'),
                          ('w', '0.15'), ('h', '0.221')):
            el = ET.SubElement(manual, qn('c:' + name))
            el.set('val', val)
        # White fill + thin navy border.
        for old in leg_el.findall(qn('c:spPr')):
            leg_el.remove(old)
        leg_spPr = ET.SubElement(leg_el, qn('c:spPr'))
        leg_fill = ET.SubElement(leg_spPr, qn('a:solidFill'))
        leg_fill_rgb = ET.SubElement(leg_fill, qn('a:srgbClr'))
        leg_fill_rgb.set('val', 'FFFFFF')
        leg_ln = ET.SubElement(leg_spPr, qn('a:ln'))
        leg_ln.set('w', '9525')
        leg_ln_fill = ET.SubElement(leg_ln, qn('a:solidFill'))
        leg_ln_rgb = ET.SubElement(leg_ln_fill, qn('a:srgbClr'))
        leg_ln_rgb.set('val', '0B2B4E')

        # 2026-05-23 third pass: takeaway moved up T 6.55 → 6.40
        # so it sits closer to the chart bottom now that the chart
        # frame ends at T 6.25 (1.58 + 4.67).
        _add_takeaway_bar(
            slide,
            "Fixed costs dominate at low Q;  the quadratic TVC overtakes at scale",
            top=Inches(6.40), fill=NAVY, text_color=WHITE,
            width=Inches(11.5), size=18,
            rounded=True, shadow=True,
        )

    s = make_diagram_slide(
        prs, page_num=55,
        section_tag=SECTION_TAG_P2,
        title="Rivian's Georgia Plant —  Cost Components",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "Decompose total cost into its components: fixed plus variable. "
        "Same Rivian cost function used on the previous slide: "
        "TC = 10,000,000 + 30,000·Q + 40·Q², so TFC = $10M (the flat "
        "gold line at 10) and TVC = 30,000·Q + 40·Q² (the rising red "
        "curve).  Their vertical sum is the navy TC curve.  The visual "
        "lets students see how the cost structure shifts as output "
        "grows: at low Q, the fixed cost dominates; at high Q, the "
        "quadratic TVC term takes over and TC tracks TVC."
    ))


def slide_56(prs):
    """Rivian's Georgia plant – Per-Unit Costs (ATC, AFC, AVC, MC).

    2026-05-23 third pass: chart redrawn to mirror original-deck
    slide 62 exactly.  Three changes vs. the prior pass:
      • Added a fourth series — AFC = 10,000,000 / Q — the
        hyperbolic average-fixed-cost curve (was missing).
      • Adopted the original deck's four-color palette: ATC blue,
        AFC orange, AVC green, MC red.
      • Replaced the navy "MC crosses ATC..." takeaway at the
        bottom with three callouts ABOVE the chart, each with a
        navy arrow pointing to the relevant feature on the curves:
          1. "ATC falls below sales price ($80k) at Q ≈ 250"
          2. "The minimum of ATC is where it crosses MC"
          3. "ATC increases again due to rising MC"

    Per-unit derivations from TC = 10M + 30k·Q + 40·Q²:
      • ATC(Q) = 10,000,000/Q + 30,000 + 40·Q     (U-shape)
      • AFC(Q) = 10,000,000/Q                      (hyperbolic)
      • AVC(Q) =                 30,000 + 40·Q     (linear up)
      • MC(Q)  =                 30,000 + 80·Q     (linear up,
                                                    twice AVC slope)

    Sales-price line is $80K.  Solving ATC(Q) = 80 gives Q = 250
    (the first crossing — ATC falling below) and Q = 1000.
    """
    def draw(slide):
        from pptx.chart.data import XyChartData
        from pptx.enum.chart import XL_MARKER_STYLE

        # ----- Per-unit formulas (matches slides 54 / 55) -------
        TFC_LOCAL = 10_000_000
        LIN_LOCAL = 30_000
        QUAD_LOCAL = 40
        atc_of = lambda q: TFC_LOCAL / q + LIN_LOCAL + QUAD_LOCAL * q
        afc_of = lambda q: TFC_LOCAL / q
        avc_of = lambda q:                  LIN_LOCAL + QUAD_LOCAL * q
        mc_of  = lambda q:                  LIN_LOCAL + 2 * QUAD_LOCAL * q

        Q_curve = list(range(100, 1001, 25))
        ATC_K = [atc_of(q) / 1000 for q in Q_curve]
        AFC_K = [afc_of(q) / 1000 for q in Q_curve]
        AVC_K = [avc_of(q) / 1000 for q in Q_curve]
        MC_K  = [mc_of(q)  / 1000 for q in Q_curve]

        chart_data = XyChartData()
        for name, vals in (('ATC', ATC_K), ('AFC', AFC_K),
                           ('AVC', AVC_K), ('MC',  MC_K)):
            s = chart_data.add_series(name)
            for q, v in zip(Q_curve, vals):
                s.add_data_point(q, v)

        # ----- Chart geometry: shifted DOWN to leave room at the
        # top for the three callout boxes (user hand-tweaked
        # 2026-05-24: T 2.10 → 2.56, H 4.95 → 4.49).
        chart_x, chart_y = Inches(2.50), Inches(2.56)
        chart_w, chart_h = Inches(8.27), Inches(4.49)
        _add_graphicframe_shadow(slide, chart_x, chart_y, chart_w, chart_h)
        gf = slide.shapes.add_chart(
            XL_CHART_TYPE.XY_SCATTER_LINES,
            chart_x, chart_y, chart_w, chart_h, chart_data,
        )
        chart = gf.chart
        chart.has_title = False

        # ----- Axes ---------------------------------------------
        xax = chart.category_axis
        xax.maximum_scale = 1000
        xax.minimum_scale = 0
        xax.major_unit = 100
        xax.has_title = True
        xax.axis_title.text_frame.text = "Q   (vehicles per week)"
        for r in xax.axis_title.text_frame.paragraphs[0].runs:
            r.font.size = Pt(18); r.font.italic = True; r.font.bold = True
            r.font.color.rgb = NAVY; r.font.name = "Calibri"
        xax.tick_labels.font.size = Pt(16)
        xax.tick_labels.font.color.rgb = NAVY
        xax.tick_labels.font.name = "Calibri"

        yax = chart.value_axis
        yax.maximum_scale = 110
        yax.minimum_scale = 0
        yax.major_unit = 10
        yax.has_title = True
        yax.axis_title.text_frame.text = "Per-unit cost   ($K)"
        for r in yax.axis_title.text_frame.paragraphs[0].runs:
            r.font.size = Pt(18); r.font.italic = True; r.font.bold = True
            r.font.color.rgb = NAVY; r.font.name = "Calibri"
        yax.tick_labels.font.size = Pt(16)
        yax.tick_labels.font.color.rgb = NAVY
        yax.tick_labels.font.name = "Calibri"

        _add_dashed_gridlines(yax._element)

        # ----- Plot area (user hand-extended to the right) ------
        # 2026-05-25: user dragged the plot area right edge to use
        # more of the chart frame; capture the resulting layout so
        # the next rebuild reproduces it.
        plot_area_el = chart._chartSpace.find('.//' + qn('c:plotArea'))
        if plot_area_el is not None:
            for old in plot_area_el.findall(qn('c:layout')):
                plot_area_el.remove(old)
            pa_layout = ET.Element(qn('c:layout'))
            pa_manual = ET.SubElement(pa_layout, qn('c:manualLayout'))
            ET.SubElement(pa_manual, qn('c:layoutTarget')).set('val', 'inner')
            ET.SubElement(pa_manual, qn('c:xMode')).set('val', 'edge')
            ET.SubElement(pa_manual, qn('c:yMode')).set('val', 'edge')
            ET.SubElement(pa_manual, qn('c:x')).set('val', '0.12872569057646513')
            ET.SubElement(pa_manual, qn('c:y')).set('val', '0.045672603842114391')
            ET.SubElement(pa_manual, qn('c:w')).set('val', '0.79662349869506943')
            ET.SubElement(pa_manual, qn('c:h')).set('val', '0.74657618660696368')
            # layout must be the first child of plotArea per schema
            plot_area_el.insert(0, pa_layout)

        # ----- Per-series styling: smooth lines, no markers -----
        # Colors mirror the original deck slide 62 palette:
        #   ATC blue, AFC orange, AVC green, MC red.
        ATC_BLUE = RGBColor(0x2E, 0x75, 0xB6)
        AFC_ORNG = RGBColor(0xED, 0x7D, 0x31)
        AVC_GRN  = RGBColor(0x70, 0xAD, 0x47)
        MC_RED   = RGBColor(0xC0, 0x00, 0x00)
        colors = [ATC_BLUE, AFC_ORNG, AVC_GRN, MC_RED]
        for i, color in enumerate(colors):
            ser = chart.plots[0].series[i]
            ser.marker.style = XL_MARKER_STYLE.NONE
            ser.format.line.color.rgb = color
            ser.format.line.width = Pt(2.75)
            ser_el = ser._element
            for old in ser_el.findall(qn('c:smooth')):
                ser_el.remove(old)
            smooth = ET.SubElement(ser_el, qn('c:smooth'))
            smooth.set('val', '1')

        # ----- Legend: vertical stacking, close together -------
        # Four entries; placed on the right side of the plot
        # area (middle vertical) where AFC has decayed and the
        # curves leave the most room.
        chart.has_legend = True
        leg = chart.legend
        leg.include_in_layout = False
        leg.font.size = Pt(18)
        leg.font.color.rgb = NAVY
        leg.font.name = "Calibri"
        leg_el = leg._element
        for old in leg_el.findall(qn('c:layout')):
            leg_el.remove(old)
        layout = ET.SubElement(leg_el, qn('c:layout'))
        manual = ET.SubElement(layout, qn('c:manualLayout'))
        for name, val in (('xMode', 'edge'), ('yMode', 'edge')):
            el = ET.SubElement(manual, qn('c:' + name))
            el.set('val', val)
        # 2026-05-24 user nudge: legend slightly up and left
        # (x 0.82 → 0.776, y 0.42 → 0.392).
        for name, val in (('x', '0.776'), ('y', '0.392'),
                          ('w', '0.14'), ('h', '0.30')):
            el = ET.SubElement(manual, qn('c:' + name))
            el.set('val', val)
        for old in leg_el.findall(qn('c:spPr')):
            leg_el.remove(old)
        leg_spPr = ET.SubElement(leg_el, qn('c:spPr'))
        leg_fill = ET.SubElement(leg_spPr, qn('a:solidFill'))
        leg_fill_rgb = ET.SubElement(leg_fill, qn('a:srgbClr'))
        leg_fill_rgb.set('val', 'FFFFFF')
        leg_ln = ET.SubElement(leg_spPr, qn('a:ln'))
        leg_ln.set('w', '9525')
        leg_ln_fill = ET.SubElement(leg_ln, qn('a:solidFill'))
        leg_ln_rgb = ET.SubElement(leg_ln_fill, qn('a:srgbClr'))
        leg_ln_rgb.set('val', '0B2B4E')

        # ----- Three callout boxes above the chart, each with
        # an arrow pointing down to a specific feature on the
        # ATC curve.  Cream rounded boxes with thin navy border —
        # same Convention-callout style used elsewhere in the
        # deck.  2026-05-24 user hand-positioned every box, every
        # arrow start/end, and bumped the text font 13 → 16 pt.
        # Arrow tips also enlarged ('med' → 'lg') so the heads
        # read clearly at the back of the room without thickening
        # the lines themselves.
        cream = RGBColor(0xFD, 0xF6, 0xE6)
        callouts = [
            # (text, box_L, box_T, box_W, box_H,
            #  arrow_start_xy, arrow_end_xy)   – all in inches
            ("ATC falls below sales price ($80k)  at Q ≈ 250",
              2.200, 1.452, 3.000, 0.705,
              (4.708, 2.212), (5.200, 3.613)),
            ("The minimum of ATC is where it crosses MC",
              5.417, 1.452, 2.850, 0.705,
              (6.666, 2.172), (6.842, 3.979)),
            ("ATC increases again due to rising MC",
              8.563, 1.508, 2.400, 0.675,
              (9.548, 2.212), (9.603, 3.743)),
        ]
        for text, bx, by, bw, bh, (sx, sy), (ex, ey) in callouts:
            _add_rounded_filled_box(
                slide,
                Inches(bx), Inches(by), Inches(bw), Inches(bh),
                label=text,
                fill=cream, text_color=NAVY, line=NAVY,
                size=16, bold=True, font="Calibri",
                corner_pct=0.08, shadow=True,
            )
            _add_arrow(slide,
                        (Inches(sx), Inches(sy)),
                        (Inches(ex), Inches(ey)),
                        color=NAVY, weight_pt=1.5,
                        head_size='lg')

    s = make_diagram_slide(
        prs, page_num=56,
        section_tag=SECTION_TAG_P2,
        title="Rivian's Georgia Plant —  Per-Unit Costs",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "Same Rivian cost function expressed per vehicle.  From "
        "TC = 10,000,000 + 30,000·Q + 40·Q²: ATC = 10M/Q + 30,000 "
        "+ 40·Q gives a U-shape (fixed-cost dilution dominates at "
        "low Q, the quadratic term takes over at high Q); AVC = "
        "30,000 + 40·Q is linear-rising; MC = 30,000 + 80·Q is "
        "linear-rising with twice the slope of AVC.  Diagnostic to "
        "highlight: MC crosses ATC exactly at the ATC minimum "
        "(here Q ≈ 500, value ≈ $70K).  The classic U-shape result "
        "from any intro micro textbook."
    ))


def slide_57(prs):
    """Problem-set preview – iPhone cost estimation.

    2026-05-25 (second pass) restructure:
      • Four cost buckets instead of three: Fixed costs (R&D,
        marketing, stores, other), Material inputs (processor is
        now part of this), Labor, and Distribution.
      • Two questions instead of one: AVC and ATC.  ATC requires
        students to compute AFC = TFC / annual worldwide iPhone
        sales (~230M / year), which is the explicit hint.
      • Retail price stated as "≈ $1,200" (was $1,199).
      • Keeps the "Problem Set Preview" badge and "you can use AI"
        tip from the first pass.
    """
    def draw(slide):
        # ----- Problem-set framing badge (top-left) -------------
        ps_badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            int(MARGIN), int(Inches(1.40)),
            int(Inches(3.10)), int(Inches(0.45)),
        )
        try: ps_badge.adjustments[0] = 0.30
        except Exception: pass
        ps_badge.fill.solid(); ps_badge.fill.fore_color.rgb = NAVY
        ps_badge.line.fill.background()
        ps_badge.shadow.inherit = False
        _add_drop_shadow(ps_badge)
        ptf = ps_badge.text_frame
        ptf.margin_left = Inches(0.10); ptf.margin_right = Inches(0.10)
        ptf.margin_top = 0; ptf.margin_bottom = 0
        ptf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = ptf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = "Problem Set Preview"
        r.font.name = "Calibri"; r.font.size = Pt(16)
        r.font.bold = True; r.font.color.rgb = WHITE

        # ----- Two core questions -------------------------------
        # 2026-05-25 (fourth pass): bumped 22 → 26 pt.
        _add_text(slide, MARGIN, Inches(1.95), RULE_W, Inches(0.55),
                  "1.  What is the AVC of a current-generation iPhone?",
                  size=26, bold=True, color=NAVY, font="Calibri",
                  align=PP_ALIGN.CENTER)
        _add_text(slide, MARGIN, Inches(2.55), RULE_W, Inches(0.55),
                  "2.  What is the ATC of a current-generation iPhone?",
                  size=26, bold=True, color=NAVY, font="Calibri",
                  align=PP_ALIGN.CENTER)

        # ----- Hint line ---------------------------------------
        # 2026-05-25 (fourth pass): dropped the parenthetical
        # "~ 200 M units" (students should research that themselves);
        # font bumped 16 → 20 pt.
        _add_text(slide, MARGIN, Inches(3.25), RULE_W, Inches(0.45),
                  "Hint:  estimate each bucket;  for AFC, divide the "
                  "fixed costs of this model by its total lifetime sales",
                  size=20, italic=True, color=GRAY, font="Calibri",
                  align=PP_ALIGN.CENTER)

        # ----- Four component cards (outlined, "$ ?" inside) ----
        # 2026-05-25 (fourth pass): card height shrunk 2.00 → 1.70
        # (taller than needed for the three text rows).  Stores
        # moved from Fixed costs → Distribution subtext.
        comps = [
            ("Fixed costs",
              "R&D, marketing,\nother"),
            ("Material inputs",
              "processor, display,\nbattery, cameras, memory, …"),
            ("Labor",
              "assembly +\nfinal test"),
            ("Distribution",
              "stores, logistics,\nshipping, channel"),
        ]
        cw = Inches(2.95)
        ch = Inches(1.70)
        gap = Inches(0.25)
        x0 = (SLIDE_W - cw * 4 - gap * 3) // 2
        for i, (name, sub) in enumerate(comps):
            card_x = x0 + (cw + gap) * i
            card_y = Inches(3.85)
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                int(card_x), int(card_y), int(cw), int(ch),
            )
            try: card.adjustments[0] = 0.08
            except Exception: pass
            card.fill.solid(); card.fill.fore_color.rgb = WHITE
            card.line.color.rgb = NAVY; card.line.width = Pt(1.5)
            card.shadow.inherit = False
            _add_drop_shadow(card)
            tf = card.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.10); tf.margin_right = Inches(0.10)
            tf.margin_top = Inches(0.08); tf.margin_bottom = Inches(0.08)
            tf.vertical_anchor = MSO_ANCHOR.TOP
            # Header (category name)
            p0 = tf.paragraphs[0]; p0.alignment = PP_ALIGN.CENTER
            r0 = p0.add_run(); r0.text = name
            r0.font.name = "Calibri"; r0.font.size = Pt(22)
            r0.font.bold = True; r0.font.color.rgb = NAVY
            # Subtext (italic gray) — bumped 12 → 16 pt
            p1 = tf.add_paragraph(); p1.alignment = PP_ALIGN.CENTER
            p1.space_before = Pt(4)
            r1 = p1.add_run(); r1.text = sub
            r1.font.name = "Calibri"; r1.font.size = Pt(16)
            r1.font.italic = True; r1.font.color.rgb = GRAY
            # "$ ?" prompt at the bottom — bumped 22 → 26 pt
            p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
            p2.space_before = Pt(6)
            r2 = p2.add_run(); r2.text = "$ ?"
            r2.font.name = "Calibri"; r2.font.size = Pt(26)
            r2.font.bold = True
            r2.font.color.rgb = RGBColor(0xC0, 0x00, 0x00)

        # ----- Anchor: retail price (the only "given") ----------
        # 2026-05-25 (fourth pass): moved down T 5.45 → 5.80 so it
        # sits clearly below the shorter cards.
        _add_text(slide, MARGIN, Inches(5.80), RULE_W, Inches(0.45),
                  "Retail price  ≈  $1,200",
                  size=22, bold=True, color=NAVY, font="Calibri",
                  align=PP_ALIGN.CENTER)

        # ----- "You can use AI" tip (gold, rounded + shadow) ----
        ai_tip = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            int((SLIDE_W - Inches(6.50)) // 2), int(Inches(6.30)),
            int(Inches(6.50)), int(Inches(0.55)),
        )
        try: ai_tip.adjustments[0] = 0.30
        except Exception: pass
        ai_tip.fill.solid(); ai_tip.fill.fore_color.rgb = GOLD
        ai_tip.line.fill.background()
        ai_tip.shadow.inherit = False
        _add_drop_shadow(ai_tip)
        atf = ai_tip.text_frame
        atf.margin_left = Inches(0.10); atf.margin_right = Inches(0.10)
        atf.margin_top = 0; atf.margin_bottom = 0
        atf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ap = atf.paragraphs[0]; ap.alignment = PP_ALIGN.CENTER
        ar = ap.add_run()
        ar.text = "Tip:  you can use AI to research the cost breakdown"
        ar.font.name = "Calibri"; ar.font.size = Pt(18)
        ar.font.bold = True; ar.font.italic = True; ar.font.color.rgb = NAVY

    s = make_diagram_slide(
        prs, page_num=57,
        section_tag=SECTION_TAG_P2,
        title="Cost Estimation:  What Does an iPhone Cost to Make?",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "Problem-set preview.  Students leave the lecture with four "
        "cost buckets to estimate (fixed costs, material inputs, "
        "labor, distribution) and a retail-price anchor of ~$1,200, "
        "then research the actual numbers on their own — AI is fine "
        "as a research aid.  Two deliverables: AVC (sum of the three "
        "per-unit buckets — material inputs, labor, distribution) "
        "AND ATC (AVC plus AFC).  Crucial framing for AFC: fixed "
        "costs (R&D, marketing, allocated stores, other) are "
        "PER MODEL, allocated over that model's TOTAL LIFETIME "
        "sales — not annual sales of all iPhones.  For a current "
        "iPhone model, lifetime worldwide sales are roughly 200M "
        "units, so AFC = TFC_model / 200M.  The pedagogical payoff "
        "is the comparison of BOTH AVC and ATC to retail price — "
        "students will see that even ATC sits well below the "
        "sticker."
    ))


def slide_58(prs):
    """Problem-set solution sketch – iPhone AVC build-up.

    2026-05-25 new content: replaces the prior in-class PollEv
    setup with a written worked-solution sketch that mirrors what
    an MBA student would hand in for the problem set previewed on
    slide 57.  Cost figures are deliberately order-of-magnitude
    teaching approximations — students should arrive at numbers
    in the same ballpark, not exact teardown estimates.

    Structure:
      • Three filled cards reproducing the slide-57 categories,
        with the bucket subtotal filled in (Processor ≈ $80,
        Other material inputs ≈ $350, Labor ≈ $30).
      • Each card lists the dominant subcomponents in small text
        beneath the bucket total.
      • Hero navy band: AVC ≈ $80 + $350 + $30 ≈ $460.
      • Retail-vs-markup comparison line (≈ $1,199 retail; ≈ $740
        gross margin per unit, ≈ 62 %).
      • Gold rounded takeaway bar restating the pedagogical
        point: build cost is a small share of retail; the rest
        covers R&D, marketing, distribution, and profit.
    """
    def draw(slide):
        RED = RGBColor(0xC0, 0x00, 0x00)
        # ----- Four filled cards (one per cost bucket) ---------
        # Card 1 (Fixed costs) is annual TOTAL — the unit on the
        # big number is "$ / yr".  Cards 2-4 are per-unit ($ /
        # iPhone).  Same visual treatment so they read as parallel
        # estimates; the explicit unit on each dollar amount makes
        # the apples-to-oranges nature obvious.
        card_w = Inches(2.95)
        card_h = Inches(2.05)
        card_gap = Inches(0.25)
        card_t = Inches(1.55)
        card_x0 = (SLIDE_W - card_w * 4 - card_gap * 3) // 2

        # 2026-05-25 (third pass): reframed Fixed-costs card from
        # annual TFC ($30B / yr) to PER-MODEL LIFETIME TFC
        # ($20B for this iPhone model's whole lifecycle), and the
        # AFC denominator from annual sales (~230M / yr) to total
        # lifetime sales (~200M units).  Card fonts bumped:
        # header 18 → 20 pt, subitems 11 → 14 pt.
        # 2026-05-25 (fourth pass): mirror slide 57's relocation of
        # "stores" from Fixed costs → Distribution.  Fixed-costs
        # subitems rebalanced so the $20B total still ties out
        # without listing "stores" separately.
        cards = [
            ("Fixed costs",
              "$20B",
              "this model (lifetime)",
              "R&D ~$10B,  marketing ~$5B,\nother ~$5B"),
            ("Material inputs",
              "$430",
              "per iPhone",
              "processor ~$80,  display ~$80,\nmemory ~$50,  battery ~$10,\ncameras ~$70,  other ~$140"),
            ("Labor",
              "$30",
              "per iPhone",
              "assembly + final test"),
            ("Distribution",
              "$20",
              "per iPhone",
              "stores, logistics,\nshipping, channel"),
        ]
        for i, (header, dollar, unit, footer) in enumerate(cards):
            card_x = card_x0 + (card_w + card_gap) * i
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                int(card_x), int(card_t), int(card_w), int(card_h),
            )
            try: card.adjustments[0] = 0.08
            except Exception: pass
            card.fill.solid(); card.fill.fore_color.rgb = WHITE
            card.line.color.rgb = NAVY; card.line.width = Pt(1.5)
            card.shadow.inherit = False
            _add_drop_shadow(card)
            tf = card.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.10); tf.margin_right = Inches(0.10)
            tf.margin_top = Inches(0.08); tf.margin_bottom = Inches(0.08)
            tf.vertical_anchor = MSO_ANCHOR.TOP
            # Header (category name)
            p0 = tf.paragraphs[0]; p0.alignment = PP_ALIGN.CENTER
            r0 = p0.add_run(); r0.text = header
            r0.font.name = "Calibri"; r0.font.size = Pt(20)
            r0.font.bold = True; r0.font.color.rgb = NAVY
            # Big dollar amount
            p1 = tf.add_paragraph(); p1.alignment = PP_ALIGN.CENTER
            p1.space_before = Pt(2)
            r1 = p1.add_run(); r1.text = dollar
            r1.font.name = "Calibri"; r1.font.size = Pt(28)
            r1.font.bold = True; r1.font.color.rgb = NAVY
            # Unit (small italic, just under the dollar amount)
            p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
            r2 = p2.add_run(); r2.text = unit
            r2.font.name = "Calibri"; r2.font.size = Pt(13)
            r2.font.italic = True; r2.font.color.rgb = GRAY
            # Footer (subcomponent list) — bumped 11 → 14 pt
            p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER
            p3.space_before = Pt(4)
            r3 = p3.add_run(); r3.text = footer
            r3.font.name = "Calibri"; r3.font.size = Pt(14)
            r3.font.italic = True; r3.font.color.rgb = GRAY

        # ----- AVC computation line (italic gray, small) -------
        _add_text(slide, MARGIN, Inches(3.75), RULE_W, Inches(0.40),
                  "AVC  =  $430  +  $30  +  $20   ≈   $480 / iPhone",
                  size=18, italic=True, bold=True, color=NAVY,
                  font="Calibri", align=PP_ALIGN.CENTER)

        # ----- AFC computation line (italic gray, small) -------
        # 2026-05-25 (third pass): denominator switched from annual
        # sales (230M / yr) to TOTAL LIFETIME sales of this model
        # (~200M units), so AFC is the per-unit allocation of the
        # model-specific lifetime TFC.
        _add_text(slide, MARGIN, Inches(4.20), RULE_W, Inches(0.40),
                  "AFC  =  $20 B  /  200 M iPhones (lifetime)   ≈   $100 / iPhone",
                  size=18, italic=True, bold=True, color=NAVY,
                  font="Calibri", align=PP_ALIGN.CENTER)

        # ----- Hero ATC band (navy filled, rounded + shadow) ---
        atc_w = Inches(9.50)
        atc_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            int((SLIDE_W - atc_w) // 2), int(Inches(4.75)),
            int(atc_w), int(Inches(0.75)),
        )
        try: atc_box.adjustments[0] = 0.20
        except Exception: pass
        atc_box.fill.solid(); atc_box.fill.fore_color.rgb = NAVY
        atc_box.line.fill.background()
        atc_box.shadow.inherit = False
        _add_drop_shadow(atc_box)
        ftf = atc_box.text_frame
        ftf.margin_left = Inches(0.15); ftf.margin_right = Inches(0.15)
        ftf.margin_top = 0; ftf.margin_bottom = 0
        ftf.vertical_anchor = MSO_ANCHOR.MIDDLE
        fp = ftf.paragraphs[0]; fp.alignment = PP_ALIGN.CENTER
        fr = fp.add_run()
        fr.text = "ATC  =  AVC + AFC  =  $480 + $100   ≈   $580 / iPhone"
        fr.font.name = "Calibri"; fr.font.size = Pt(24)
        fr.font.bold = True; fr.font.color.rgb = WHITE

        # ----- Comparison schematic: AVC | AFC | ATC | Retail --
        pill_data = [
            ("AVC\n$480",    NAVY,                          WHITE),
            ("+ AFC\n$100",  GOLD,                          NAVY),
            ("= ATC\n$580",  NAVY,                          WHITE),
            ("Retail\n$1,200", RGBColor(0xFD, 0xF6, 0xE6),  RED),  # cream w/ red text
        ]
        pill_w = Inches(2.30)
        pill_h = Inches(0.75)
        pill_gap = Inches(0.15)
        n = len(pill_data)
        pill_total_w = pill_w * n + pill_gap * (n - 1)
        pill_x0 = (SLIDE_W - pill_total_w) // 2
        pill_y = Inches(5.65)
        for i, (text, fill, txtcol) in enumerate(pill_data):
            pill_x = pill_x0 + (pill_w + pill_gap) * i
            pill = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                int(pill_x), int(pill_y), int(pill_w), int(pill_h),
            )
            try: pill.adjustments[0] = 0.25
            except Exception: pass
            pill.fill.solid(); pill.fill.fore_color.rgb = fill
            if fill == RGBColor(0xFD, 0xF6, 0xE6):
                pill.line.color.rgb = RED; pill.line.width = Pt(1.5)
            else:
                pill.line.fill.background()
            pill.shadow.inherit = False
            _add_drop_shadow(pill)
            ptf2 = pill.text_frame
            ptf2.margin_left = Inches(0.05); ptf2.margin_right = Inches(0.05)
            ptf2.margin_top = Inches(0.03); ptf2.margin_bottom = Inches(0.03)
            ptf2.vertical_anchor = MSO_ANCHOR.MIDDLE
            ptf2.word_wrap = True
            # Multi-line: line1 label, line2 dollar.  Split on '\n'.
            lines = text.split('\n')
            for li, line in enumerate(lines):
                if li == 0:
                    pp = ptf2.paragraphs[0]
                else:
                    pp = ptf2.add_paragraph()
                pp.alignment = PP_ALIGN.CENTER
                rr = pp.add_run(); rr.text = line
                rr.font.name = "Calibri"
                rr.font.size = Pt(14 if li == 0 else 18)
                rr.font.bold = True
                rr.font.color.rgb = txtcol

        # ----- Gold takeaway bar (rounded + shadow) -----------
        _add_takeaway_bar(
            slide,
            "Even ATC sits well below retail  —  the rest is gross "
            "profit margin on each iPhone",
            top=Inches(6.55), fill=GOLD, text_color=NAVY,
            width=Inches(11.5), size=18,
            rounded=True, shadow=True,
        )

    s = make_diagram_slide(
        prs, page_num=58,
        section_tag=SECTION_TAG_P2,
        title="Solution Sketch:  Estimating the iPhone Cost Stack",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "Walk through the kind of answer an MBA student would hand "
        "in for the problem set previewed last slide.  Four buckets, "
        "order-of-magnitude estimates that basic research lands on:  "
        "(1) Fixed costs FOR THIS MODEL OVER ITS LIFETIME ~$20B — "
        "R&D allocated to this generation ~$10B, marketing ~$5B, "
        "other corporate overhead ~$5B.  (2) Material inputs ~$430 "
        "per iPhone — A-series chip ~$80 plus display ~$80, memory "
        "~$50, battery ~$10, multi-lens cameras ~$70, and the long "
        "tail of modem, antennas, housing, packaging ~$140.  "
        "(3) Labor (assembly + final test) ~$30 per iPhone.  "
        "(4) Distribution (Apple Retail stores, logistics, "
        "shipping, channel) ~$20 per iPhone — stores treated as a "
        "distribution-channel cost amortised per unit.  Sum the "
        "three per-unit buckets → AVC ≈ $480 / iPhone.  Critical "
        "framing for AFC: divide the model's lifetime TFC by its "
        "TOTAL lifetime worldwide sales (~200M units) — NOT by "
        "annual sales.  AFC ≈ $20B / 200M = $100 / iPhone.  "
        "ATC = AVC + AFC ≈ $580 / iPhone.  Retail ~$1,200 minus "
        "ATC ~$580 = ~$620 gross profit margin per iPhone "
        "(~52%), consistent with Apple's reported iPhone gross "
        "margin in the mid-40% to low-50% range.  Pedagogical "
        "payoff: even after loading in fixed costs, ATC is "
        "roughly half the sticker price — the rest is profit."
    ))


def slide_59(prs):
    """Solution: AVC of iPhone 17 ≈ $580 (vs. $1,199 retail)."""
    def draw(slide):
        # Teardown image on the left
        _add_source_image(slide, 60, "rId4",
                          left=Inches(0.5), top=Inches(1.85),
                          height=Inches(4.6))

        # Numbers on the right
        _add_text(slide, Inches(7.5), Inches(2.0), Inches(5.4), Inches(0.55),
                  "Retail price:",
                  size=22, color=GRAY, font="Calibri")
        _add_text(slide, Inches(7.5), Inches(2.55), Inches(5.4), Inches(0.6),
                  "$1,199",
                  size=36, bold=True, color=NAVY, font="Calibri")
        _add_text(slide, Inches(7.5), Inches(3.4), Inches(5.4), Inches(0.55),
                  "Total variable cost  (TVC):",
                  size=22, color=GRAY, font="Calibri")
        _add_text(slide, Inches(7.5), Inches(3.95), Inches(5.4), Inches(0.6),
                  "≈  $580",
                  size=36, bold=True, color=GOLD, font="Calibri")

        # Missing components (small list)
        _add_text(slide, Inches(7.5), Inches(4.85), Inches(5.4), Inches(0.4),
                  "Plus missing components:",
                  size=14, italic=True, color=GRAY, font="Calibri")
        miss = ["Shipping & handling", "Customer service", "Warranty costs"]
        for i, m in enumerate(miss):
            _add_text(slide, Inches(7.7), Inches(5.2 + i * 0.32),
                      Inches(5.2), Inches(0.3),
                      f"–  {m}",
                      size=14, color=GRAY, font="Calibri")

        _add_takeaway_bar(
            slide,
            "About half of retail goes to fixed-cost recovery, R&D, retail margin",
            top=Inches(6.5), fill=NAVY, width=Inches(11.0),
        )

    s = make_diagram_slide(
        prs, page_num=59,
        section_tag=SECTION_TAG_P2,
        title="AVC of iPhone 17  ≈  $580",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "Reveal: AVC ≈ $580 for iPhone 17 (vs. ~$1,199 retail). About "
        "half of retail. The rest – fixed-cost recovery, gross margin – "
        "funds Apple's R&D, retail network, and ecosystem. Students "
        "consistently overestimate this number."
    ))


def slide_60(prs):
    """iPhone naïve cost function – total cost (line chart, hand-drawn)."""
    def draw(slide):
        # Axes
        AX_L = Inches(1.6)
        AX_T = Inches(2.0)
        AX_W = Inches(7.5)
        AX_H = Inches(4.2)

        # Y-axis line
        _add_arrow(slide,
                   start_xy=(int(AX_L), int(AX_T + AX_H)),
                   end_xy=(int(AX_L), int(AX_T)),
                   color=NAVY, weight_pt=1.5, head=True)
        # X-axis line
        _add_arrow(slide,
                   start_xy=(int(AX_L), int(AX_T + AX_H)),
                   end_xy=(int(AX_L + AX_W), int(AX_T + AX_H)),
                   color=NAVY, weight_pt=1.5, head=True)

        # Axis labels
        _add_text(slide, AX_L - Inches(1.0), AX_T - Inches(0.45),
                   Inches(2.0), Inches(0.4),
                   "$  (Total cost)", size=18, bold=True, color=NAVY,
                   font="Calibri")
        _add_text(slide, AX_L + AX_W - Inches(2.5),
                   AX_T + AX_H + Inches(0.15),
                   Inches(3.5), Inches(0.4),
                   "Q  (quantity produced)", size=18, bold=True, color=NAVY,
                   font="Calibri")

        # Fixed-cost intercept (horizontal mark) at ~ TFC level
        tfc_y = AX_T + AX_H - Inches(0.9)   # TFC ~ 20% of vertical
        _add_text(slide, AX_L - Inches(0.7), tfc_y - Inches(0.2),
                   Inches(0.7), Inches(0.4),
                   "TFC", size=18, bold=True, color=GOLD, font="Calibri")
        # Short tick from y-axis
        _add_arrow(slide,
                   start_xy=(int(AX_L - Inches(0.05)), int(tfc_y)),
                   end_xy=(int(AX_L + Inches(0.05)), int(tfc_y)),
                   color=GOLD, weight_pt=2.0, head=False)

        # The TC line: starts at (AX_L, tfc_y), goes up-right
        line_end_x = AX_L + AX_W - Inches(0.5)
        line_end_y = AX_T + Inches(0.4)
        _add_arrow(slide,
                   start_xy=(int(AX_L), int(tfc_y)),
                   end_xy=(int(line_end_x), int(line_end_y)),
                   color=NAVY, weight_pt=3.0, head=False)

        # Equation callout on the right (OMML – upright acronyms + italic Q)
        eq_xml = (
            _omml_text('TC') +
            _omml_text(' = ') +
            _omml_text('TFC') +
            _omml_text(' + ') +
            _omml_text('500') +
            _omml_text(' · ') +
            _omml_run('Q')
        )
        _add_math_equation(
            slide, Inches(9.5), Inches(3.0), Inches(3.5), Inches(0.95),
            eq_xml, size_pt=24, color=NAVY, fill=GOLD,
        )

        # Slope annotation
        _add_callout_box(
            slide, Inches(5.8), Inches(4.3), Inches(2.4), Inches(0.45),
            "slope  =  $500",
            fill=NAVY, text_color=WHITE, size=14, bold=True,
        )

        _add_takeaway_bar(
            slide,
            "Linear naïve TC :  fixed cost  +  constant marginal cost  ($500 / unit)",
            top=Inches(6.55), fill=NAVY, width=Inches(11.5),
        )

    s = make_diagram_slide(
        prs, page_num=60,
        section_tag=SECTION_TAG_P2,
        title="Naïve Linear Cost Function:  Total Cost View",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "iPhone naïve cost function – total cost as output rises. "
        "Linear: a fixed cost plus a constant marginal cost. This is "
        "the simplest model of a cost function; we'll see more complex "
        "shapes later (capacity constraints, increasing returns)."
    ))


def slide_61(prs):
    """iPhone naïve cost function – per-unit (constant MC line)."""
    def draw(slide):
        AX_L = Inches(1.6)
        AX_T = Inches(2.0)
        AX_W = Inches(7.5)
        AX_H = Inches(4.2)

        # Y-axis
        _add_arrow(slide,
                   start_xy=(int(AX_L), int(AX_T + AX_H)),
                   end_xy=(int(AX_L), int(AX_T)),
                   color=NAVY, weight_pt=1.5, head=True)
        # X-axis
        _add_arrow(slide,
                   start_xy=(int(AX_L), int(AX_T + AX_H)),
                   end_xy=(int(AX_L + AX_W), int(AX_T + AX_H)),
                   color=NAVY, weight_pt=1.5, head=True)

        _add_text(slide, AX_L - Inches(1.0), AX_T - Inches(0.45),
                   Inches(2.0), Inches(0.4),
                   "$ per unit", size=18, bold=True, color=NAVY,
                   font="Calibri")
        _add_text(slide, AX_L + AX_W - Inches(2.5),
                   AX_T + AX_H + Inches(0.15),
                   Inches(3.5), Inches(0.4),
                   "Q  (quantity produced)", size=18, bold=True, color=NAVY,
                   font="Calibri")

        # Constant horizontal line at $500 mark
        mc_y = AX_T + AX_H - Inches(2.0)   # half-way up
        _add_arrow(slide,
                   start_xy=(int(AX_L + Inches(0.2)), int(mc_y)),
                   end_xy=(int(AX_L + AX_W - Inches(0.3)), int(mc_y)),
                   color=NAVY, weight_pt=3.0, head=False)
        # Y-axis tick label "500"
        _add_text(slide, AX_L - Inches(0.7), mc_y - Inches(0.2),
                   Inches(0.7), Inches(0.4),
                   "500", size=18, bold=True, color=GOLD, font="Calibri",
                   align=PP_ALIGN.RIGHT)

        # Line label (OMML)
        eq_label = (
            _omml_text('MC') +
            _omml_text(' = ') +
            _omml_text('AVC') +
            _omml_text(' = ') +
            _omml_text('$500')
        )
        _add_math_equation(
            slide,
            AX_L + AX_W // 2 - Inches(1.8),
            mc_y - Inches(0.65),
            Inches(3.6), Inches(0.55),
            eq_label, size_pt=22, color=NAVY,
        )

        # Annotation callout (OMML inside)
        callout_xml = (
            _omml_text('When ') +
            _omml_text('MC') +
            _omml_text(' is constant, it equals ') +
            _omml_text('AVC')
        )
        _add_math_equation(
            slide,
            Inches(9.5), Inches(3.0), Inches(3.5), Inches(1.5),
            callout_xml, size_pt=18, color=NAVY, fill=GOLD,
        )

        _add_takeaway_bar(
            slide,
            "Per-unit view :  constant marginal cost,  AC falls as fixed cost spreads",
            top=Inches(6.55), fill=GOLD, text_color=NAVY,
            width=Inches(11.5),
        )

    s = make_diagram_slide(
        prs, page_num=61,
        section_tag=SECTION_TAG_P2,
        title="Naïve Linear Cost Function:  Per-Unit View",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "Same data, per unit. Note the constant MC and (implicit) "
        "declining AC – the classic shape when fixed costs spread over "
        "more units. With constant MC, AVC = MC at every Q. "
        "AC = AVC + AFC, so AC starts above MC and falls toward it as "
        "Q → ∞."
    ))


# --------------------------------------------------------------------------
# Batch 5 – §2.2 Long-Run Costs & Economies of Scale  (slides 63-74)
# --------------------------------------------------------------------------


def slide_62(prs):
    """Bridge slide: real-world cost functions can be non-linear (U-shape
    MC). We will keep MC linear where possible – the iPhone toy model is
    fine for most decisions; the Rivian quadratic was the exception."""
    def draw(slide):
        # Two source images side by side: TC/VC/FC curves + U-shape MC
        _add_source_image(slide, 63, "rId1",
                          left=Inches(0.5), top=Inches(2.0),
                          height=Inches(3.6))
        _add_source_image(slide, 63, "rId2",
                          left=Inches(7.0), top=Inches(2.0),
                          height=Inches(3.6))
        # Captions
        _add_text(slide, Inches(0.5), Inches(5.65),
                  Inches(6.0), Inches(0.3),
                  "Total / variable / fixed cost — convex at high Q",
                  size=14, italic=True, color=GRAY,
                  font="Calibri", align=PP_ALIGN.CENTER)
        _add_text(slide, Inches(7.0), Inches(5.65),
                  Inches(6.0), Inches(0.3),
                  "Marginal cost — U-shape (high at low + high Q)",
                  size=14, italic=True, color=GRAY,
                  font="Calibri", align=PP_ALIGN.CENTER)

        _add_takeaway_bar(
            slide,
            "Keep MC linear when possible — use the U-shape only when scale really matters",
            top=Inches(6.5), fill=GOLD, text_color=NAVY,
            width=Inches(12.0),
        )

    s = make_diagram_slide(
        prs, page_num=62,
        section_tag=SECTION_TAG_P2,
        title="More Complex Cost Functions:  When MC Isn't Linear",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "Reality check: cost functions can be more complex than a "
        "straight line. The classic textbook shape is a U for marginal "
        "cost – MC high at low Q (under-utilised plant), falls as you "
        "hit the sweet spot, then rises again at high Q (over-stretched "
        "capacity, overtime, congestion). We saw the convex piece in the "
        "Rivian quadratic. For most decisions in this course we will "
        "still use linear MC – it captures the right intuition without "
        "the algebra. Pull out the U-shape only when scale effects are "
        "the whole point."
    ))


def slide_63(prs):
    """Section divider – Part 2.2: Long-Run Costs & Economies of Scale.
    Mirror of slide_30 (Part 1.2) using the Part-2 highlight."""
    s = make_section_agenda(
        prs, page_num=63,
        current_part_idx=1,
        section_tag=SECTION_TAG_DIV,
        title="Part 2.2:  Long-Run Costs & Economies of Scale",
    )
    _set_notes(s, (
        "Transitioning from the cost concepts (fixed/variable/marginal/"
        "average) to what happens in the LONG run, when the plant size "
        "itself can change. Three big ideas next: long-run vs short-run "
        "cost curves, the LR-AC envelope, and economies of scale."
    ))


def slide_64(prs):
    """Short-run v. long-run costs – two-column comparison."""
    def draw(slide):
        col_w = Inches(6.0)
        col_h = Inches(0.85)
        gap = Inches(0.2)
        x_l = MARGIN
        x_r = MARGIN + col_w + gap
        y0 = Inches(2.0)

        # Headers
        _add_filled_box(slide, x_l, y0, col_w, col_h,
                         "Short Run",
                         fill=NAVY, text_color=WHITE, size=26, bold=True)
        _add_filled_box(slide, x_r, y0, col_w, col_h,
                         "Long Run",
                         fill=NAVY, text_color=WHITE, size=26, bold=True)

        # Body bullets
        sr_items = [
            "Capital  (K)  is FIXED",
            "Labor  (L)  is flexible",
            "Plant size already chosen",
            "→  Cost of changing Q given the plant you have",
        ]
        lr_items = [
            "BOTH K and L flexible",
            "Choose plant size from scratch",
            "Pick the optimal input mix for each Q",
            "→  Lowest cost of producing any Q",
        ]

        for i, (l, r) in enumerate(zip(sr_items, lr_items)):
            yi = y0 + col_h + Inches(0.25) + Inches(0.7) * i
            _add_text(slide, x_l + Inches(0.1), yi, col_w - Inches(0.2),
                      Inches(0.6), l,
                      size=18, color=NAVY, font="Calibri")
            _add_text(slide, x_r + Inches(0.1), yi, col_w - Inches(0.2),
                      Inches(0.6), r,
                      size=18, color=NAVY, font="Calibri")

        # Inequality at the bottom of the body
        eq_xml = (
            _omml_text('TC') +
            _omml_sub(_omml_text(''), _omml_text('SR')) +
            _omml_text('  ≥  ') +
            _omml_text('TC') +
            _omml_sub(_omml_text(''), _omml_text('LR'))
        )
        _add_math_equation(
            slide,
            (SLIDE_W - Inches(7.0)) // 2, Inches(5.55),
            Inches(7.0), Inches(0.7),
            eq_xml, size_pt=26, color=WHITE, fill=NAVY,
        )

        _add_takeaway_bar(
            slide,
            "Long-run costs are the lower envelope:  more freedom  ⇒  weakly cheaper",
            top=Inches(6.5), fill=GOLD, text_color=NAVY,
            width=Inches(11.0),
        )

    s = make_diagram_slide(
        prs, page_num=64,
        section_tag=SECTION_TAG_P2_LR,
        title="Short-Run vs. Long-Run Costs",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "Key distinction. Short run: you're stuck with the plant you "
        "have, so you're paying for capacity you may not be using (or "
        "scrambling to get more). Long run: you get to build the right "
        "plant for the output you want. With more freedom comes weakly "
        "lower cost – long-run TC is never higher than short-run TC at "
        "any given Q, because the long run is the option to pick the "
        "best short-run plant for that Q."
    ))


def slide_65(prs):
    """LR-AC envelope schematic.

    A clean schematic: three plant-size SAC bands across a quantity axis,
    with the LR-AC line traced as their lower envelope. Drawn with shapes
    rather than curves (which python-pptx doesn't render smoothly).
    """
    def draw(slide):
        # Plot bounding box
        AX_L = Inches(1.2)
        AX_R = Inches(12.5)
        AX_T = Inches(2.1)
        AX_B = Inches(5.8)
        AX_W = AX_R - AX_L
        AX_H = AX_B - AX_T

        # Axes
        _add_arrow(slide, (AX_L, AX_B), (AX_R, AX_B),
                   color=NAVY, weight_pt=1.5, head=True)
        _add_arrow(slide, (AX_L, AX_B), (AX_L, AX_T),
                   color=NAVY, weight_pt=1.5, head=True)
        _add_text(slide, AX_R - Inches(0.9), AX_B + Inches(0.05),
                  Inches(1.0), Inches(0.3),
                  "Q  (output)", size=14, italic=True, color=GRAY,
                  font="Calibri")
        _add_text(slide, AX_L - Inches(0.9), AX_T - Inches(0.05),
                  Inches(1.0), Inches(0.3),
                  "$ / unit", size=14, italic=True, color=GRAY,
                  font="Calibri", align=PP_ALIGN.RIGHT)

        # Three plant-size U-shape SAC bands.  Each is drawn as a smooth
        # custGeom curve (two cubic Béziers meeting at the minimum) so the
        # parabolic feel of the original deck comes through, rather than
        # piecewise straight segments.
        def draw_sac(label, x_min, x_max, y_min, y_left, y_right,
                      color=GRAY):
            x_mid = (x_min + x_max) // 2
            bb_left = x_min
            bb_top = min(y_left, y_right)
            bb_right = x_max
            bb_bottom = y_min
            bb_w = bb_right - bb_left
            bb_h = bb_bottom - bb_top
            if bb_w <= 0 or bb_h <= 0:
                return
            # Convert to path-coord space (100000 × 100000)
            def px(real_x):
                return int(100000 * (real_x - bb_left) / bb_w)
            def py(real_y):
                return int(100000 * (real_y - bb_top) / bb_h)

            # Cubic-Bezier control points for a smooth U:
            #   - tangent at the endpoints is steeply VERTICAL (down on the
            #     left side, up on the right side) — produces the
            #     parabolic descent into / out of the trough;
            #   - tangent at the minimum is HORIZONTAL — flat bottom of U.
            # Previous placement (CP1 horizontal from P0) collapsed the
            # curve to an L-shape; this placement gives a proper U.
            P0L = (px(x_min),  py(y_left))
            PMD = (px(x_mid),  py(y_min))
            P3R = (px(x_max),  py(y_right))
            seg1_cp1 = (P0L[0] + (PMD[0] - P0L[0]) // 10,
                         P0L[1] + (PMD[1] - P0L[1]) * 7 // 10)
            seg1_cp2 = (PMD[0] - (PMD[0] - P0L[0]) * 3 // 10, PMD[1])
            seg2_cp1 = (PMD[0] + (P3R[0] - PMD[0]) * 3 // 10, PMD[1])
            seg2_cp2 = (P3R[0] - (P3R[0] - PMD[0]) // 10,
                         P3R[1] + (PMD[1] - P3R[1]) * 7 // 10)

            r_hex = f'{color[0]:02X}{color[1]:02X}{color[2]:02X}'
            custgeom_xml = (
                f'<a:custGeom xmlns:a="{A_NS}">'
                f'<a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/>'
                f'<a:rect l="0" t="0" r="100000" b="100000"/>'
                f'<a:pathLst>'
                f'<a:path w="100000" h="100000" fill="none" stroke="1">'
                f'<a:moveTo><a:pt x="{px(x_min)}" y="{py(y_left)}"/></a:moveTo>'
                f'<a:cubicBezTo>'
                f'<a:pt x="{seg1_cp1[0]}" y="{seg1_cp1[1]}"/>'
                f'<a:pt x="{seg1_cp2[0]}" y="{seg1_cp2[1]}"/>'
                f'<a:pt x="{px(x_mid)}" y="{py(y_min)}"/>'
                f'</a:cubicBezTo>'
                f'<a:cubicBezTo>'
                f'<a:pt x="{seg2_cp1[0]}" y="{seg2_cp1[1]}"/>'
                f'<a:pt x="{seg2_cp2[0]}" y="{seg2_cp2[1]}"/>'
                f'<a:pt x="{px(x_max)}" y="{py(y_right)}"/>'
                f'</a:cubicBezTo>'
                f'</a:path></a:pathLst>'
                f'</a:custGeom>'
            )

            shp = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                int(bb_left), int(bb_top), int(bb_w), int(bb_h),
            )
            shp.fill.background()
            shp.line.color.rgb = color
            shp.line.width = Pt(2.25)
            shp.shadow.inherit = False
            spPr = shp._element.spPr
            for old in spPr.findall(qn('a:prstGeom')):
                spPr.remove(old)
            custgeom = ET.fromstring(custgeom_xml)
            xfrm = spPr.find(qn('a:xfrm'))
            if xfrm is not None:
                xfrm.addnext(custgeom)
            else:
                spPr.insert(0, custgeom)

            # Label above the minimum
            _add_text(slide,
                       x_mid - Inches(0.9),
                       y_min - Inches(0.45),
                       Inches(1.8), Inches(0.3),
                       label, size=12, italic=True, color=color,
                       font="Calibri", align=PP_ALIGN.CENTER)

        # Three SAC bands placed along x — minima fall as you go right
        # (economies of scale).
        w = AX_W
        b1_x_min = AX_L + int(0.05 * w); b1_x_max = AX_L + int(0.32 * w)
        b1_y_min = AX_T + int(0.55 * AX_H)
        b1_y_left  = AX_T + int(0.15 * AX_H)
        b1_y_right = AX_T + int(0.20 * AX_H)
        draw_sac("SAC  (small plant)", b1_x_min, b1_x_max,
                  b1_y_min, b1_y_left, b1_y_right)

        b2_x_min = AX_L + int(0.30 * w); b2_x_max = AX_L + int(0.62 * w)
        b2_y_min = AX_T + int(0.65 * AX_H)
        b2_y_left  = AX_T + int(0.25 * AX_H)
        b2_y_right = AX_T + int(0.30 * AX_H)
        draw_sac("SAC  (medium plant)", b2_x_min, b2_x_max,
                  b2_y_min, b2_y_left, b2_y_right)

        b3_x_min = AX_L + int(0.60 * w); b3_x_max = AX_L + int(0.95 * w)
        b3_y_min = AX_T + int(0.75 * AX_H)
        b3_y_left  = AX_T + int(0.35 * AX_H)
        b3_y_right = AX_T + int(0.40 * AX_H)
        draw_sac("SAC  (large plant)", b3_x_min, b3_x_max,
                  b3_y_min, b3_y_left, b3_y_right)

        # Lower-envelope LAC — passes EXACTLY through each SAC minimum
        # so the three U-curves visibly touch the envelope at their lowest
        # point.  Two straight segments connect the three (mid_x, y_min)
        # points;  pedagogically this is the textbook envelope rendering.
        b1_x_mid = (b1_x_min + b1_x_max) // 2
        b2_x_mid = (b2_x_min + b2_x_max) // 2
        b3_x_mid = (b3_x_min + b3_x_max) // 2
        _add_arrow(slide,
                    (b1_x_mid, b1_y_min),
                    (b2_x_mid, b2_y_min),
                    color=GOLD, weight_pt=3.0, head=False)
        _add_arrow(slide,
                    (b2_x_mid, b2_y_min),
                    (b3_x_mid, b3_y_min),
                    color=GOLD, weight_pt=3.0, head=False)
        _add_text(slide, AX_R - Inches(2.6), b3_y_min - Inches(0.55),
                  Inches(2.5), Inches(0.3),
                  "LAC  =  lower envelope",
                  size=14, bold=True, italic=True, color=GOLD,
                  font="Calibri")

        _add_takeaway_bar(
            slide,
            "LR-AC is the lower envelope of all SAC curves:  pick the best plant for each Q",
            top=Inches(6.4), fill=GOLD, text_color=NAVY,
            width=Inches(12.0),
        )

    s = make_diagram_slide(
        prs, page_num=65,
        section_tag=SECTION_TAG_P2_LR,
        title="LR Average Cost is the Lower Envelope of SR Curves",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "Classic envelope diagram. Each plant size has its own short-run "
        "AC curve – U-shaped, with a minimum at the output it was "
        "designed for. The long-run AC is the lower envelope of all "
        "these short-run curves: for each output level Q, you pick the "
        "plant size that produces it at the lowest AC. In this drawing, "
        "the LAC slopes down – meaning economies of scale: bigger plant "
        "→ lower AC at its respective optimum."
    ))


def slide_66(prs):
    """Economies of scale: definition + the three cases."""
    def draw(slide):
        # Question header
        _add_text(slide, MARGIN, Inches(1.85), RULE_W, Inches(0.55),
                  "What happens to long-run AC as output grows?",
                  size=24, italic=True, color=GRAY,
                  font="Calibri", align=PP_ALIGN.CENTER)

        # Three cases as horizontal cards
        case_w = Inches(4.0)
        case_h = Inches(2.0)
        gap = Inches(0.25)
        x0 = (SLIDE_W - case_w * 3 - gap * 2) // 2
        y0 = Inches(2.8)

        _add_filled_box(slide, x0, y0, case_w, case_h,
                         "Economies of Scale\n\nLAC FALLS with Q\n\n(bigger ⇒ cheaper / unit)",
                         fill=NAVY, text_color=WHITE,
                         size=18, bold=True)
        _add_filled_box(slide, x0 + (case_w + gap), y0, case_w, case_h,
                         "Constant Returns\n\nLAC is FLAT in Q\n\n(size doesn't matter)",
                         fill=GRAY, text_color=WHITE,
                         size=18, bold=True)
        _add_filled_box(slide, x0 + 2 * (case_w + gap), y0, case_w, case_h,
                         "Diseconomies of Scale\n\nLAC RISES with Q\n\n(too big to manage)",
                         fill=NAVY, text_color=WHITE,
                         size=18, bold=True)

        # Why bullet header
        _add_text(slide, MARGIN, Inches(5.05), RULE_W, Inches(0.4),
                  "Why?  Two big drivers of economies of scale:",
                  size=18, bold=True, color=NAVY,
                  font="Calibri", align=PP_ALIGN.CENTER)
        _add_text(slide, MARGIN, Inches(5.50), RULE_W, Inches(0.4),
                  "(1) input prices fall as you grow      (2) technology favours larger scale (increasing returns)",
                  size=16, italic=True, color=GRAY,
                  font="Calibri", align=PP_ALIGN.CENTER)

        _add_takeaway_bar(
            slide,
            "Economies of scale is a COST concept;  returns to scale is a TECHNOLOGY concept",
            top=Inches(6.4), fill=GOLD, text_color=NAVY,
            width=Inches(12.0),
        )

    s = make_diagram_slide(
        prs, page_num=66,
        section_tag=SECTION_TAG_P2_LR,
        title="Economies of Scale:  Three Possible Patterns",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "Economies of scale describes what happens to long-run average "
        "cost as output rises. Three cases: falling (EoS), flat "
        "(constant), rising (DEoS). Drivers are either input-price "
        "effects – you get bulk discounts, or you bid up prices when "
        "you're huge – or pure technology effects. Important: "
        "'economies of scale' is a COST-side concept, while 'returns to "
        "scale' (a textbook term) is a TECHNOLOGY concept. The two "
        "overlap but aren't the same. A firm can have constant returns "
        "to scale but enjoy economies of scale from bulk pricing."
    ))


def slide_67(prs):
    """Technological reasons for economies of scale."""
    bullets = [
        ("Specialisation and division of labor", 0),
        ("E.g., Ford assembly line vs. one craftsman building the whole car", 1),
        ("Lumpiness / indivisibilities of inputs", 0),
        ("R&D, network infrastructure, brand investment", 1),
        ("Same fixed asset serves more customers as scale grows", 1),
        ("Geometry — volume scales faster than surface", 0),
        ("Cargo ship, aircraft fuselage, pipeline:  capacity grows with r²", 1),
        ("Doubling materials more than doubles useful capacity", 1),
        ("Appropriate technology shifts with scale", 0),
        ("Big firms can run dedicated lines, automation, AI/data infra", 1),
    ]
    s = make_content_bulleted(
        prs, page_num=67,
        section_tag=SECTION_TAG_P2_LR,
        title="Technological Reasons for Economies of Scale",
        bullets=bullets,
        size=24, sub_size=22, line_spacing_pts=10,
    )
    _set_notes(s, (
        "Why does technology often favour larger scale? Four classic "
        "drivers. Specialisation – Adam Smith's pin factory – workers "
        "get better at narrow tasks. Lumpiness – you can't build half a "
        "network or half an R&D lab, so the fixed cost is the same "
        "whether you have 1M or 100M users. Geometry – for cylinders "
        "and tanks, volume grows faster than the surface area you have "
        "to build, so cost per unit of capacity drops with size. And "
        "scale unlocks DIFFERENT technologies entirely: only Amazon-"
        "scale firms can justify their own AI infrastructure or "
        "fulfilment robotics."
    ))


def slide_68(prs):
    """Embraer ERJ-145 vs. Boeing 787 – scale economies in aviation."""
    def draw(slide):
        # Two airplane cards side by side
        card_w = Inches(6.0)
        card_h = Inches(4.4)
        gap = Inches(0.3)
        x_l = (SLIDE_W - card_w * 2 - gap) // 2
        x_r = x_l + card_w + gap
        y0 = Inches(1.95)

        # Left card – Embraer
        _add_outlined_box(slide, x_l, y0, card_w, card_h,
                          "", fill=WHITE, line=NAVY, line_w=1.5)
        _add_text(slide, x_l, y0 + Inches(0.1), card_w, Inches(0.4),
                  "Embraer ERJ-145  ·  Regional Jet",
                  size=18, bold=True, color=NAVY,
                  font="Calibri", align=PP_ALIGN.CENTER)
        _add_source_image(slide, 69, "rId1",
                          left=x_l + Inches(0.3), top=y0 + Inches(0.55),
                          width=card_w - Inches(0.6))
        # Stats
        stats_y = y0 + Inches(2.4)
        stat_lines = [
            "List price:    ~ $25 M",
            "Seats:           50 passengers",
            "Cost / flight-hour:    ~ $1,400",
            "Cost / passenger-hour:   ≈ $28",
        ]
        for i, line in enumerate(stat_lines):
            _add_text(slide, x_l + Inches(0.4),
                      stats_y + Inches(0.4) * i,
                      card_w - Inches(0.8), Inches(0.35),
                      line, size=16, color=NAVY, font="Calibri")

        # Right card – Boeing 787
        _add_outlined_box(slide, x_r, y0, card_w, card_h,
                          "", fill=WHITE, line=NAVY, line_w=1.5)
        _add_text(slide, x_r, y0 + Inches(0.1), card_w, Inches(0.4),
                  "Boeing 787-9  ·  Wide-Body",
                  size=18, bold=True, color=NAVY,
                  font="Calibri", align=PP_ALIGN.CENTER)
        _add_source_image(slide, 69, "rId2",
                          left=x_r + Inches(0.3), top=y0 + Inches(0.55),
                          width=card_w - Inches(0.6))
        stat_lines_r = [
            "List price:    ~ $290 M",
            "Seats:           ~ 290 passengers",
            "Cost / flight-hour:    ~ $9,000",
            "Cost / passenger-hour:   ≈ $31",
        ]
        for i, line in enumerate(stat_lines_r):
            _add_text(slide, x_r + Inches(0.4),
                      stats_y + Inches(0.4) * i,
                      card_w - Inches(0.8), Inches(0.35),
                      line, size=16, color=NAVY, font="Calibri")

        _add_takeaway_bar(
            slide,
            "Bigger plane  →  similar (or lower) cost per passenger-hour:  geometry + load",
            top=Inches(6.45), fill=GOLD, text_color=NAVY,
            width=Inches(12.0),
        )

    s = make_diagram_slide(
        prs, page_num=68,
        section_tag=SECTION_TAG_P2_LR,
        title="Economies of Scale in Aviation:  ERJ-145 vs. 787",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "Concrete example. Both aircraft burn fuel, pay pilots, pay "
        "landing fees – many of those costs scale with the airframe, "
        "not with passenger count. Spread across far more seats, the "
        "big plane achieves a similar or lower cost per passenger-hour. "
        "(Numbers are illustrative, based on FAA operating-cost data.) "
        "This is why long-haul routes use wide-body and short regional "
        "hops use small jets: the economics of scale depend on the load."
    ))


def slide_69(prs):
    """Reasons for diseconomies of scale."""
    bullets = [
        ("Coordination, communication, control  get harder", 0),
        ("More layers of management between strategy and the line", 1),
        ("Information has to travel up and down a longer chain", 1),
        ("Monitoring a bigger workforce is disproportionately costly", 0),
        ("Misaligned incentives, free-riding, principal-agent problems", 1),
        ("Bureaucracy and staff functions grow super-linearly", 0),
        ("HR, accounting, legal, compliance hire to support scale", 1),
        ("Real-world signals", 0),
        ("Boeing's recent quality issues at scale  (2024-25)", 1),
        ("Big-tech reorgs to break ranks into smaller, accountable units", 1),
    ]
    s = make_content_bulleted(
        prs, page_num=69,
        section_tag=SECTION_TAG_P2_LR,
        title="Reasons for Diseconomies of Scale",
        bullets=bullets,
        size=24, sub_size=22, line_spacing_pts=10,
    )
    _set_notes(s, (
        "Why bigger isn't always cheaper. Coordination costs explode "
        "with size: longer reporting chains, more meetings, more "
        "alignment overhead. Monitoring a 50-person team is one thing; "
        "monitoring 5,000 is qualitatively different. Bureaucracy "
        "scales super-linearly – legal, HR, compliance all expand "
        "faster than headcount. Recent business-press examples drive "
        "this home: Boeing's quality issues at scale, and the wave of "
        "big-tech reorganisations explicitly framed around 'getting "
        "back to small-team velocity'."
    ))


def slide_70(prs):
    """Economies of scope – producing 2+ products together."""
    def draw(slide):
        # Definition box at top
        _add_filled_box(slide, MARGIN, Inches(1.95),
                         RULE_W, Inches(0.9),
                         "Economies of scope:   producing 2+ related products together is cheaper than separately",
                         fill=NAVY, text_color=WHITE,
                         size=18, bold=True)

        # Sub-bullet drivers (left side)
        bullets = [
            ("Shared input production", 0),
            ("Shared engineering know-how & R&D", 0),
            ("Shared brand, sales channel, marketing", 0),
            ("Shared supply chain & customer data", 0),
        ]
        _add_hierarchical_bullets(
            slide,
            left=MARGIN, top=Inches(3.15),
            width=Inches(6.5), height=Inches(2.8),
            items=bullets,
            size=20, line_spacing_pts=10,
        )

        # Example image on right – Airbus A380 + A318
        _add_source_image(slide, 71, "rId1",
                          left=Inches(7.5), top=Inches(3.15),
                          width=Inches(5.5))
        _add_text(slide, Inches(7.5), Inches(5.85),
                  Inches(5.5), Inches(0.25),
                  "British Airways A380 + A318  (CC BY-SA, Wikimedia)",
                  size=11, italic=True, color=GRAY,
                  font="Calibri", align=PP_ALIGN.CENTER)
        _add_text(slide, Inches(7.5), Inches(3.0),
                  Inches(5.5), Inches(0.25),
                  "Example: Airbus's A380 and A318 share engineering, supply chain, brand",
                  size=12, italic=True, color=NAVY,
                  font="Calibri", align=PP_ALIGN.CENTER)

        _add_takeaway_bar(
            slide,
            "Scope ≠ scale:  one firm, MANY products  →  cheaper than splitting them up",
            top=Inches(6.5), fill=GOLD, text_color=NAVY,
            width=Inches(11.5),
        )

    s = make_diagram_slide(
        prs, page_num=70,
        section_tag=SECTION_TAG_P2_LR,
        title="Economies of Scope:  Cheaper Together than Apart",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "Scope is the cousin of scale – but about the BREADTH of "
        "products a firm makes, not the depth of any one. Same firm, "
        "many products, sharing capabilities that are expensive to "
        "build: engineering knowledge, supply chain, brand. Airbus "
        "designed the A380 and A318 with shared engineering DNA – "
        "cockpit commonalities, pilot training, supplier base. Apple "
        "does it across iPhone/iPad/Mac (shared silicon, OS, retail). "
        "Recognising scope is what separates a firm that 'diversifies' "
        "into unrelated junk from one that genuinely lowers costs."
    ))


def slide_71(prs):
    """Amazon – scale, scope, or both? Discussion."""
    def draw(slide):
        _add_text(slide, MARGIN, Inches(1.95), RULE_W, Inches(0.55),
                  "Amazon is the textbook case for BOTH at once  —  where exactly?",
                  size=22, italic=True, bold=True, color=NAVY,
                  font="Calibri", align=PP_ALIGN.CENTER)

        # Two columns: Scale | Scope
        col_w = Inches(6.0)
        col_h = Inches(0.75)
        gap = Inches(0.3)
        x_l = (SLIDE_W - col_w * 2 - gap) // 2
        x_r = x_l + col_w + gap
        y_hdr = Inches(2.7)

        _add_filled_box(slide, x_l, y_hdr, col_w, col_h,
                         "Economies of SCALE",
                         fill=NAVY, text_color=WHITE,
                         size=22, bold=True)
        _add_filled_box(slide, x_r, y_hdr, col_w, col_h,
                         "Economies of SCOPE",
                         fill=NAVY, text_color=WHITE,
                         size=22, bold=True)

        scale_items = [
            "AWS — massive data-center fixed costs spread",
            "Fulfilment network density (FBA)",
            "Bargaining power with suppliers",
        ]
        scope_items = [
            "Prime  =  shipping + video + music + grocery",
            "Customer data shared across retail/ads/AWS",
            "Devices (Alexa, Kindle) lever the brand",
        ]
        y_items = y_hdr + col_h + Inches(0.25)
        for i, (l, r) in enumerate(zip(scale_items, scope_items)):
            _add_text(slide, x_l + Inches(0.2),
                      y_items + Inches(0.55) * i,
                      col_w - Inches(0.4), Inches(0.5),
                      "•  " + l,
                      size=17, color=NAVY, font="Calibri")
            _add_text(slide, x_r + Inches(0.2),
                      y_items + Inches(0.55) * i,
                      col_w - Inches(0.4), Inches(0.5),
                      "•  " + r,
                      size=17, color=NAVY, font="Calibri")

        _add_discussion_break(slide, width=Inches(5.0))

    s = make_diagram_slide(
        prs, page_num=71,
        section_tag=SECTION_TAG_P2_LR,
        title="Amazon:  Economies of Scale,  Scope,  or Both?",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "Discussion prompt. Amazon is the canonical 'both' case. "
        "Scale: AWS is a fixed-cost-heavy business where capacity is "
        "shared across millions of customers. Fulfilment is denser the "
        "more orders flow through. Suppliers offer better terms at "
        "scale. Scope: Prime bundles shipping + video + music + "
        "grocery; advertising, devices, and AWS all benefit from "
        "customer data created in retail. The strategic question for "
        "students: which of these is most defensible, and which is "
        "just incremental?"
    ))


def slide_72(prs):
    """Shark Tank mini-case setup (video + group discussion)."""
    def draw(slide):
        # Setup line
        _add_text(slide, MARGIN, Inches(1.95), RULE_W, Inches(0.5),
                  "Watch first:   vimeo.com/236977187   (focus 4:50 – 5:40)",
                  size=20, italic=True, bold=True, color=NAVY,
                  font="Calibri", align=PP_ALIGN.CENTER)

        # Two PollEV question cards
        card_w = Inches(6.0)
        card_h = Inches(2.0)
        gap = Inches(0.3)
        x_l = (SLIDE_W - card_w * 2 - gap) // 2
        x_r = x_l + card_w + gap
        y0 = Inches(2.8)

        _add_outlined_box(slide, x_l, y0, card_w, card_h,
                          "PollEV  ·  Q1\n\nAre there economies of scale\nin this business?",
                          fill=WHITE, line=NAVY, text_color=NAVY,
                          size=20, bold=True, line_w=2.0)
        _add_outlined_box(slide, x_r, y0, card_w, card_h,
                          "PollEV  ·  Q2\n\nWhich deal would you choose?\n"
                          "$100K, royalty 25¢/can   vs.   $75K, 15% equity",
                          fill=WHITE, line=NAVY, text_color=NAVY,
                          size=18, bold=True, line_w=2.0)

        # Cue bullet under the cards
        _add_text(slide, MARGIN, Inches(5.1), RULE_W, Inches(0.5),
                  "Look for:  volume last year vs. this year,  and average cost per can",
                  size=18, italic=True, color=GRAY,
                  font="Calibri", align=PP_ALIGN.CENTER)
        _add_text(slide, MARGIN, Inches(5.55), RULE_W, Inches(0.5),
                  "Then estimate profit per can to compare the two deals on an apples-to-apples basis",
                  size=18, italic=True, color=GRAY,
                  font="Calibri", align=PP_ALIGN.CENTER)

        _add_discussion_break(slide, width=Inches(5.0))

    s = make_diagram_slide(
        prs, page_num=72,
        section_tag=SECTION_TAG_P2_LR,
        title="Mini-Case:  Shark Tank Pitch — Group Discussion",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "Group discussion. Watch the pitch (Vimeo link), focus 4:50 – "
        "5:40 where they discuss sales volumes and costs. Two questions "
        "on PollEV. First: is this a business with economies of scale? "
        "Compare last year's volume and unit cost with this year's "
        "projection. Second: of the two deals the sharks offered, which "
        "is more founder-friendly? Comparison requires estimating "
        "profit per can."
    ))


def slide_73(prs):
    """Shark Tank mini-case – the numbers / solution."""
    def draw(slide):
        # Two columns: Volume/Cost evidence  |  Deal comparison
        col_w = Inches(6.0)
        gap = Inches(0.3)
        x_l = (SLIDE_W - col_w * 2 - gap) // 2
        x_r = x_l + col_w + gap
        y0 = Inches(1.95)

        # Left – volume + AC
        _add_filled_box(slide, x_l, y0, col_w, Inches(0.7),
                         "Volume & average cost",
                         fill=NAVY, text_color=WHITE,
                         size=20, bold=True)
        evid = [
            "Sales last year:     135K cans",
            "Sales this year:    300K cans   ↑",
            "Avg. cost last yr:   $1.30 / can",
            "Avg. cost this yr:   $1.10 / can   ↓",
            "⇒  AC falls as volume rises  =  Economies of Scale",
        ]
        for i, line in enumerate(evid):
            _add_text(slide, x_l + Inches(0.2),
                      y0 + Inches(0.85) + Inches(0.4) * i,
                      col_w - Inches(0.4), Inches(0.35),
                      line, size=16,
                      color=GOLD if line.startswith("⇒") else NAVY,
                      bold=line.startswith("⇒"),
                      font="Calibri")

        # Right – deal comparison
        _add_filled_box(slide, x_r, y0, col_w, Inches(0.7),
                         "Deals on a per-can basis",
                         fill=NAVY, text_color=WHITE,
                         size=20, bold=True)
        deals = [
            "Wholesale price:    $2.69 / can",
            "Profit per can:     $2.69 − $1.10  ≈  $1.50",
            "Deal A:  25¢ royalty / can  =  16.6% of profit",
            "Deal B:  15% equity",
            "⇒  Very similar in expected NPV terms",
        ]
        for i, line in enumerate(deals):
            _add_text(slide, x_r + Inches(0.2),
                      y0 + Inches(0.85) + Inches(0.4) * i,
                      col_w - Inches(0.4), Inches(0.35),
                      line, size=16,
                      color=GOLD if line.startswith("⇒") else NAVY,
                      bold=line.startswith("⇒"),
                      font="Calibri")

        _add_takeaway_bar(
            slide,
            "EoS is real here.  Royalty and equity deals look similar  —  pick based on control & risk",
            top=Inches(6.45), fill=GOLD, text_color=NAVY,
            width=Inches(12.5),
        )

    s = make_diagram_slide(
        prs, page_num=73,
        section_tag=SECTION_TAG_P2_LR,
        title="Shark Tank Solution:  Scale + Deal Comparison",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "Reveal. The numbers show a clear AC drop – $1.30 → $1.10 per "
        "can – as volume more than doubles. That's economies of scale "
        "in the can-production business. Wholesale price $2.69 leaves "
        "about $1.50 profit per can. Deal A (25¢ royalty) is roughly "
        "16.6% of profit per can; Deal B (15% equity) is structurally "
        "similar in expected value. The real decision then turns on "
        "non-cash factors: control, signalling, dilution path, and the "
        "founder's view on probability of upside."
    ))


# --------------------------------------------------------------------------
# Backup section (slides 75–77) — bare cover slide plus two placeholder
# slides linked from slide 16's "Very high / Very low MPL image"
# annotations.  2026-05-19 (user request).
# --------------------------------------------------------------------------

SECTION_TAG_BACKUP = "Module 3 · Backup"


def _add_back_button(slide, *, fill=None, top=None):
    """Lower-right block-arrow 'back' button.  Tagged so the hyperlink-
    wiring pass at the end of build_deck() can find it.

    ``fill`` defaults to light gray (#D9D9D9); pass an explicit
    ``RGBColor`` to override.  ``top`` defaults to Inches(6.50).
    """
    if fill is None:
        fill = RGBColor(0xD9, 0xD9, 0xD9)
    btn_w = Inches(1.30)
    btn_h = Inches(0.55)
    btn_x = SLIDE_W - MARGIN - btn_w
    btn_y = top if top is not None else Inches(6.50)
    shp = _add_arrow_shape(slide, btn_x, btn_y, btn_w, btn_h,
                            direction="left", fill=fill)
    _add_drop_shadow(shp)
    tf = shp.text_frame
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "back"
    r.font.name = "Calibri"
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = NAVY
    # Tag the shape so the post-build wiring pass can identify it.
    nvSpPr = shp._element.find(qn('p:nvSpPr'))
    cNvPr = nvSpPr.find(qn('p:cNvPr')) if nvSpPr is not None else None
    if cNvPr is not None:
        cNvPr.set('name', 'BackButton')
    return shp


def slide_75_backup_cover(prs):
    """Bare cover slide — large 'BACKUP' word centred.  Footer kept for
    page-number navigation; top bar omitted for visual quiet."""
    slide = _blank_slide(prs)
    _add_text(slide, MARGIN, Inches(2.7), RULE_W, Inches(2.0),
               "BACKUP",
               size=140, bold=True, color=NAVY, font="Calibri",
               align=PP_ALIGN.CENTER)
    # 2026-05-20: was 75 — shifted down one when slide 37 was deleted.
    _draw_footer(slide, FOOTER_TEXT, 74)
    return slide


def _backup_chrome(slide, *, title, background_image=None,
                    title_left, title_top, title_width, title_height,
                    back_button_top=None):
    """Minimalist chrome for the MPL-illustration backup slides.

    2026-05-19 (manual, latest): only three elements remain per slide —
    the full-slide background image, the cream title pill on the right,
    and the light-gray back button.  Top bar / section tag, page number,
    rule lines, and gold accents are all dropped.
    """
    # 1. Background image (whole-slide; LARGE PNG ~2.5 MB each).
    if background_image is not None and background_image.exists():
        slide.shapes.add_picture(
            str(background_image), 0, 0,
            width=SLIDE_W, height=SLIDE_H,
        )
    # 2. Title pill on a cream banner.  Geometry passed by the caller
    #    so each slide can fine-tune (slide 77 places its title flush
    #    against the top edge with a taller two-line-ready box).
    title_box = _add_text(
        slide, title_left, title_top, title_width, title_height,
        title,
        size=32, bold=True, color=NAVY, font="Calibri",
        align=PP_ALIGN.LEFT,
    )
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = RGBColor(0xEE, 0xEC, 0xE1)
    # 3. Light-gray back button, bottom-right.  Wired to slide 16 by
    #    the post-build hyperlink pass in build_deck().
    _add_back_button(slide, top=back_button_top)


def slide_76_backup_high_mpl(prs):
    """Backup slide — 'Very High MPL in the Rivian Plant' (linked from slide 16).

    2026-05-19 (manual, latest): top bar removed, page number removed —
    only the background image, the cream title pill, and the light-gray
    back button remain.  Title position nudged right (6.31" → 6.495").
    """
    slide = _blank_slide(prs)
    _backup_chrome(
        slide,
        # 2026-05-19 (manual, later): collapsed the 4-space gap between
        # "MPL" and "in" down to a single space, so the title reads as
        # one phrase rather than a label + subtitle.
        title="Very High MPL in the Rivian Plant",
        background_image=OUT_DIR / "Background Material" / "Module 3 - Rivian Plant -- Empty.png",
        # 2026-05-19 (manual): title pill pinned to the right — left
        # edge dragged inward (6.425 → 6.870, w 6.748 → 6.303) so the
        # right edge stays at 13.173 against the slide's right margin.
        title_left=Inches(6.870), title_top=Inches(0.180),
        title_width=Inches(6.303), title_height=Inches(0.539),
    )
    _set_notes(slide, (
        "Backup slide showing what a very high marginal product of labor "
        "looks like in practice — the start of Rivian's hiring ramp, where "
        "each additional worker adds a lot of output.  Linked from the "
        "'Very high MPL image' annotation on slide 16; the bottom-right "
        "back button returns to slide 16."
    ))


def slide_77_backup_low_mpl(prs):
    """Backup slide — 'Very Low MPL in the Rivian Plant' (linked from slide 16).

    2026-05-19 (manual, latest): same minimalist treatment as slide 76,
    but the title sits flush against the top of the slide (top=0) in a
    taller box, and the back button is hand-shifted down (~0.27").
    """
    slide = _blank_slide(prs)
    _backup_chrome(
        slide,
        # 2026-05-19 (manual): wording extended to flag the
        # near-zero-or-negative tail of the MPL curve.
        title="Very Low (or negative) MPL in the Rivian Plant",
        background_image=OUT_DIR / "Background Material" / "Module 3 - Rivian Plant -- Crowded.png",
        # Title pill moved right + narrowed (6.877 → 8.190, w 6.456 →
        # 5.143) so it tucks against the slide's right edge.
        title_left=Inches(8.190), title_top=Inches(0.0),
        title_width=Inches(5.143), title_height=Inches(1.077),
        back_button_top=Inches(6.77),
    )
    _set_notes(slide, (
        "Backup slide showing what a very low marginal product of labor "
        "looks like in practice — the saturated tail of Rivian's hiring "
        "curve, where each additional worker adds almost nothing.  Linked "
        "from the 'Very low MPL image' annotation on slide 16; the "
        "bottom-right back button returns to slide 16."
    ))


# --------------------------------------------------------------------------
# Layout-stripping surgery (kept from previous version)
# --------------------------------------------------------------------------

KEEP_LAYOUT = 'slideLayout7.xml'
LAYOUT_DISPLAY_NAME = '405 Slides Layout'


def strip_unused_layouts(pptx_path: Path):
    src = pptx_path
    tmp = pptx_path.with_suffix(pptx_path.suffix + '.tmp')

    with zipfile.ZipFile(src, 'r') as zin:
        names = zin.namelist()

        layouts_to_drop = []
        for n in names:
            if n.startswith('ppt/slideLayouts/') and n.endswith('.xml'):
                fname = n.rsplit('/', 1)[-1]
                if fname != KEEP_LAYOUT:
                    layouts_to_drop.append(n)
            elif n.startswith('ppt/slideLayouts/_rels/') and n.endswith('.xml.rels'):
                fname = n.rsplit('/', 1)[-1].replace('.rels', '')
                if fname != KEEP_LAYOUT:
                    layouts_to_drop.append(n)

        drop_set = set(layouts_to_drop)

        master_rels_xml = zin.read('ppt/slideMasters/_rels/slideMaster1.xml.rels').decode('utf-8')
        rels_root = ET.fromstring(master_rels_xml.encode('utf-8'))
        REL_NS = 'http://schemas.openxmlformats.org/package/2006/relationships'

        dropped_rids = []
        for rel in rels_root.findall(f'{{{REL_NS}}}Relationship'):
            target = rel.get('Target', '')
            if 'slideLayouts/' in target:
                if not target.endswith(KEEP_LAYOUT):
                    dropped_rids.append(rel.get('Id'))
                    rels_root.remove(rel)

        new_master_rels = ET.tostring(
            rels_root, xml_declaration=True, encoding='UTF-8',
            standalone=True,
        ).decode('utf-8')

        master_xml = zin.read('ppt/slideMasters/slideMaster1.xml').decode('utf-8')
        for rid in dropped_rids:
            master_xml = re.sub(
                rf'<p:sldLayoutId\s+id="\d+"\s+r:id="{rid}"\s*/>',
                '',
                master_xml,
            )

        ct_xml = zin.read('[Content_Types].xml').decode('utf-8')
        for n in layouts_to_drop:
            if n.endswith('.xml'):
                part_name = '/' + n
                ct_xml = re.sub(
                    rf'<Override\s+PartName="{re.escape(part_name)}"\s+ContentType="[^"]*"\s*/>',
                    '',
                    ct_xml,
                )

        kept_layout_xml = zin.read(f'ppt/slideLayouts/{KEEP_LAYOUT}').decode('utf-8')
        kept_layout_xml = re.sub(
            r'<p:cSld\s+name="[^"]*"',
            f'<p:cSld name="{LAYOUT_DISPLAY_NAME}"',
            kept_layout_xml,
            count=1,
        )
        if f'name="{LAYOUT_DISPLAY_NAME}"' not in kept_layout_xml:
            kept_layout_xml = kept_layout_xml.replace(
                '<p:cSld>', f'<p:cSld name="{LAYOUT_DISPLAY_NAME}">', 1,
            )

        with zipfile.ZipFile(tmp, 'w', zipfile.ZIP_DEFLATED) as zout:
            for n in names:
                if n in drop_set:
                    continue
                if n == 'ppt/slideMasters/_rels/slideMaster1.xml.rels':
                    zout.writestr(n, new_master_rels)
                elif n == 'ppt/slideMasters/slideMaster1.xml':
                    zout.writestr(n, master_xml)
                elif n == '[Content_Types].xml':
                    zout.writestr(n, ct_xml)
                elif n == f'ppt/slideLayouts/{KEEP_LAYOUT}':
                    zout.writestr(n, kept_layout_xml)
                else:
                    zout.writestr(n, zin.read(n))

    shutil.move(str(tmp), str(src))


def build_deck(output_name="Module 3_clean.pptx"):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Front matter
    slide_1(prs)
    slide_2(prs)
    slide_announcements(prs)     # page 3 — midterm logistics (reintroduced 2026-05-15)
    slide_3(prs)                 # page 4 onwards
    slide_4(prs)
    slide_5(prs)
    # Big-picture concept map – at page 6 (replaces old textual outline)
    slide_concept_map(prs)

    # Part 1 §1.1 Short Run
    slide_7(prs)
    slide_8(prs)
    slide_9(prs)
    slide_10(prs)
    slide_11(prs)
    slide_short_run_agenda(prs)  # page 13 — NEW intra-Part-1 divider (2026-05-19)
    slide_12(prs)                # page 14 — Short Run MPL concept intro
    slide_mpl_data(prs)          # page 15 — MPL data (matches original slide 17)
    slide_13(prs)                # page 16 onwards
    slide_14(prs)
    slide_15(prs)
    slide_16(prs)
    # slide_17(prs)  — MERGED into slide_16; function kept for reference only
    slide_18(prs)
    slide_19(prs)
    slide_20(prs)
    # slide_21(prs)  — MERGED into slide_22; function kept for reference
    slide_22(prs)
    slide_22b(prs)               # page 23 — NEW: numerical MRPL solution

    # Part 1 §1.1b Wage Searchers
    slide_23(prs)
    # slide_24(prs)  — MERGED into slide_23 (2026-05-18); function kept for reference
    slide_25(prs)
    slide_26(prs)
    slide_27(prs)
    slide_28(prs)
    slide_29(prs)

    # Part 1.2 Long Run (with section divider)
    slide_30(prs)
    slide_31(prs)
    slide_32(prs)
    # slide_33(prs) — 2026-05-19: deleted by user; the headline rule now
    # lives on slide_32's page, with the three sub-bullets rescued there.
    slide_34(prs)
    slide_35(prs)
    slide_36(prs)
    slide_37(prs)
    slide_38(prs)
    slide_39(prs)
    slide_40(prs)

    # Part 2 section divider
    slide_41(prs)

    # Part 2 §2.1 Cost Concepts
    slide_42(prs)
    slide_43(prs)
    slide_44(prs)
    slide_45(prs)
    slide_46(prs)
    slide_47(prs)
    slide_48(prs)
    slide_49(prs)
    slide_50(prs)
    slide_51(prs)
    slide_52(prs)
    slide_53(prs)
    slide_54(prs)
    slide_55(prs)
    slide_56(prs)
    slide_57(prs)
    slide_58(prs)
    slide_59(prs)
    slide_60(prs)
    slide_61(prs)

    # Part 2 §2.2 Long-Run Costs & Economies of Scale
    slide_62(prs)
    slide_63(prs)
    slide_64(prs)
    slide_65(prs)
    slide_66(prs)
    slide_67(prs)
    slide_68(prs)
    slide_69(prs)
    slide_70(prs)
    slide_71(prs)
    slide_72(prs)
    slide_73(prs)

    # Backup section (slides 75-77).  Slide 16's right-side annotations
    # ("Very high MPL image" / "Very low MPL image") link to the
    # appropriate backup slide; each backup slide's bottom-right "back"
    # button returns to slide 16.
    slide_75_backup_cover(prs)
    slide_76_backup_high_mpl(prs)
    slide_77_backup_low_mpl(prs)

    # Wire slide-jump hyperlinks now that all targets exist.
    slides = list(prs.slides)
    s16 = slides[15]   # 0-indexed: deck slide 16
    s76 = slides[74]
    s77 = slides[75]
    # 1. Slide-16 annotation textboxes → backup slides (matched by text).
    for shape in s16.shapes:
        if not shape.has_text_frame:
            continue
        body = shape.text_frame.text
        if "Very high MPL" in body:
            tgt = s76
        elif "Very low MPL" in body:
            tgt = s77
        else:
            continue
        # Only hyperlink the run that holds the actual label text;
        # the leading "➤  " prefix run stays plain (no underline,
        # no hyperlink).
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                if "MPL image" in r.text:
                    _add_slide_jump_hyperlink_run(s16, r, tgt)
    # 2. Back buttons on backup slides → slide 16.
    for src in (s76, s77):
        for shape in src.shapes:
            if shape.name == 'BackButton':
                _add_slide_jump_hyperlink_shape(src, shape, s16)

    out = OUT_DIR / output_name
    prs.save(out)
    strip_unused_layouts(out)
    return out


if __name__ == "__main__":
    import sys
    name = sys.argv[1] if len(sys.argv) > 1 else "Module 3_clean.pptx"
    out = build_deck(name)
    print(f"Wrote {out}")
