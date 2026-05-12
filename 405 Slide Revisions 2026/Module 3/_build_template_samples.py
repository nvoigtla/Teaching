"""
Build the 405 Slides Template – 6-layout reference deck.

Outputs "405 Slides Template.pptx" with one sample slide for each of the
six layouts the Module 3 rebuild needs:

    1. Title slide       (cover-style, no navy top bar)
    2. Section header    (you-are-here agenda: current Part in navy,
                          other Parts in light grey)
    3. Content bulleted  (action title + real PPT bullets in navy,
                          navy body text, properly aligned)
    4. Content 2-column  (concept | definition; the original sample)
    5. Poll slide        (action title = question, A-D auto-num bullets
                          in navy, no QR box)
    6. Closing synthesis (top half = recap, bottom half = preview,
                          gold rule splits the body)

Visual system: navy #0B2B4E, gold #E09F3E, light-gray rule #C8CDD3,
body gray #555B66, faded grey for inactive agenda items #B0B5BC.
Calibri throughout. 0.7 cm side margins for all main horizontal
elements. Two-gold-segment rhythm (under title and above footer) on
content slides for parallel visual pacing.

Built from scratch with python-pptx using only fonts that ship with
Windows / Office. No external assets, no attribution required.
"""

from pathlib import Path

from lxml import etree as ET
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

OUT_DIR = Path(__file__).parent

# Standard 16:9 slide
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Visual constants
NAVY = RGBColor(0x0B, 0x2B, 0x4E)
GRAY = RGBColor(0x55, 0x5B, 0x66)
FADED = RGBColor(0xB0, 0xB5, 0xBC)        # faded grey for inactive agenda
RULE = RGBColor(0xC8, 0xCD, 0xD3)
GOLD = RGBColor(0xE0, 0x9F, 0x3E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# 0.7 cm side margin
MARGIN = Inches(0.7 / 2.54)               # ≈ 0.276"
RULE_W = SLIDE_W - MARGIN * 2
GOLD_W = Inches(2.2)


# --------------------------------------------------------------------------
# Low-level helpers
# --------------------------------------------------------------------------

def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _add_text(slide, left, top, width, height, text, *,
              size=18, bold=False, color=RGBColor(0, 0, 0),
              font="Calibri", align=PP_ALIGN.LEFT, italic=False,
              anchor=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    if anchor is not None:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def _add_rect(slide, left, top, width, height, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


# --------------------------------------------------------------------------
# Native PowerPoint bullet support via direct OOXML manipulation
# (python-pptx does not expose this; we set the paragraph properties
# pPr children — buClr, buSzPct, buFont, buChar or buAutoNum — in the
# schema-required order.)
# --------------------------------------------------------------------------

_PPR_BULLET_TAGS = ('buClr', 'buSzTx', 'buSzPct', 'buSzPts',
                    'buFontTx', 'buFont', 'buNone', 'buAutoNum',
                    'buChar', 'buBlip')


def _clear_bullet(pPr):
    for tag in _PPR_BULLET_TAGS:
        for el in pPr.findall(qn(f'a:{tag}')):
            pPr.remove(el)


def _set_bullet_char(paragraph, *, char='▪', color=NAVY, font='Calibri',
                     mar_l=342900, indent=-342900, size_pct=100):
    """Attach a real PPT character bullet to *paragraph*.

    mar_l   – left margin where wrapped text aligns, in EMU (914400 = 1")
    indent  – bullet hang from mar_l (negative = bullet sits to the left)
    """
    p = paragraph._p
    pPr = p.get_or_add_pPr()
    pPr.set('marL', str(mar_l))
    pPr.set('indent', str(indent))
    _clear_bullet(pPr)

    if color is not None:
        buClr = ET.SubElement(pPr, qn('a:buClr'))
        srgb = ET.SubElement(buClr, qn('a:srgbClr'))
        srgb.set('val', '{:02X}{:02X}{:02X}'.format(color[0], color[1], color[2]))
    if size_pct and size_pct != 100:
        buSzPct = ET.SubElement(pPr, qn('a:buSzPct'))
        buSzPct.set('val', str(size_pct * 1000))
    buFont = ET.SubElement(pPr, qn('a:buFont'))
    buFont.set('typeface', font)
    buChar = ET.SubElement(pPr, qn('a:buChar'))
    buChar.set('char', char)


def _set_bullet_autonum(paragraph, *, scheme='alphaUcPeriod', color=NAVY,
                        font='Calibri', mar_l=457200, indent=-457200,
                        start_at=None):
    """Attach a PPT auto-numbered bullet (A. B. C. by default)."""
    p = paragraph._p
    pPr = p.get_or_add_pPr()
    pPr.set('marL', str(mar_l))
    pPr.set('indent', str(indent))
    _clear_bullet(pPr)

    if color is not None:
        buClr = ET.SubElement(pPr, qn('a:buClr'))
        srgb = ET.SubElement(buClr, qn('a:srgbClr'))
        srgb.set('val', '{:02X}{:02X}{:02X}'.format(color[0], color[1], color[2]))
    buFont = ET.SubElement(pPr, qn('a:buFont'))
    buFont.set('typeface', font)
    buAutoNum = ET.SubElement(pPr, qn('a:buAutoNum'))
    buAutoNum.set('type', scheme)
    if start_at is not None:
        buAutoNum.set('startAt', str(start_at))


def _add_bulleted_list(slide, left, top, width, height, items, *,
                       size=22, color=NAVY, bold=False, font='Calibri',
                       bullet_color=NAVY, bullet_char='▪',
                       bullet_size_pct=100,
                       line_spacing_pts=12,
                       autonum_scheme=None):
    """Render *items* as a real-bulleted list in a single text frame.

    If autonum_scheme is given (e.g. 'alphaUcPeriod') the bullets become
    auto-numbered (A. B. C.) instead of character bullets.
    """
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    for i, text in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        # Spacing before each item except the first
        if i > 0:
            pPr = p._p.get_or_add_pPr()
            spcBef = ET.SubElement(pPr, qn('a:spcBef'))
            pts = ET.SubElement(spcBef, qn('a:spcPts'))
            pts.set('val', str(line_spacing_pts * 100))

        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color

        if autonum_scheme:
            _set_bullet_autonum(p, scheme=autonum_scheme, color=bullet_color,
                                font=font)
        else:
            _set_bullet_char(p, char=bullet_char, color=bullet_color,
                             font=font, size_pct=bullet_size_pct)
    return box


# --------------------------------------------------------------------------
# Shared chrome (top bar + footer rhythm) for content-family layouts
# --------------------------------------------------------------------------

def _draw_top_bar(slide, section_tag):
    bar_h = Inches(0.34)
    _add_rect(slide, 0, 0, SLIDE_W, bar_h, NAVY)
    _add_text(slide, MARGIN, 0, Inches(12), bar_h,
              section_tag.upper(), size=11, bold=True,
              color=WHITE, font="Calibri",
              anchor=MSO_ANCHOR.MIDDLE)


def _draw_action_title(slide, title, gold_len=GOLD_W):
    _add_text(slide, MARGIN, Inches(0.55), RULE_W, Inches(0.7),
              title, size=32, bold=True, color=NAVY, font="Calibri")
    _add_rect(slide, MARGIN, Inches(1.25), RULE_W, Inches(0.02), RULE)
    _add_rect(slide, MARGIN, Inches(1.235), gold_len, Inches(0.05), GOLD)


def _draw_footer(slide, footer_text, page_num):
    _add_rect(slide, 0, Inches(7.15), SLIDE_W, Inches(0.02), RULE)
    _add_rect(slide, MARGIN, Inches(7.135), GOLD_W, Inches(0.05), GOLD)
    _add_text(slide, MARGIN, Inches(7.22), Inches(11), Inches(0.3),
              footer_text, size=10, color=GRAY)
    _add_text(slide, Inches(12.6), Inches(7.22), Inches(0.5), Inches(0.3),
              str(page_num), size=10, color=GRAY, align=PP_ALIGN.RIGHT)


FOOTER_TEXT = "Management 405  ·  Module 3  ·  Production and Costs"


def _draw_poll_pill(slide):
    """Small navy 'POLL' pill, top-right, sitting just below the top bar.

    Replaces the heavier "?" badge. Echoes the navy of the top bar so it
    reads as part of the same chrome, not as an alert. Visible as the
    on-slide cue that this is a poll/exercise.
    """
    pill_w = Inches(1.05)
    pill_h = Inches(0.34)
    pill_x = SLIDE_W - MARGIN - pill_w
    pill_y = Inches(0.55)

    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  pill_x, pill_y, pill_w, pill_h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = NAVY
    shp.line.fill.background()
    shp.shadow.inherit = False
    # Tighten corner radius for a sleeker pill
    try:
        shp.adjustments[0] = 0.5
    except Exception:
        pass

    _add_text(slide, pill_x, pill_y, pill_w, pill_h,
              "POLL", size=12, bold=True, color=WHITE, font="Calibri",
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Small gold accent dot to the left of the pill, vertical-centered
    dot_size = Inches(0.14)
    dot_x = pill_x - Inches(0.16) - dot_size
    dot_y = pill_y + (pill_h - dot_size) // 2
    _add_rect(slide, dot_x, dot_y, dot_size, dot_size, GOLD)


# --------------------------------------------------------------------------
# Layout 1 – Title slide (cover-style, no navy top bar)
# --------------------------------------------------------------------------

def slide_title(prs):
    slide = _blank_slide(prs)

    # Module title – navy, very large
    _add_text(slide, MARGIN, Inches(2.0), RULE_W, Inches(1.3),
              "Production and Costs",
              size=60, bold=True, color=NAVY, font="Calibri",
              align=PP_ALIGN.CENTER)

    # Module label – larger per 2026-05-11 session-2 feedback
    _add_text(slide, MARGIN, Inches(3.35), RULE_W, Inches(0.75),
              "Module 3",
              size=40, bold=True, color=GOLD, font="Calibri",
              align=PP_ALIGN.CENTER)

    # Gold accent bar
    accent_w = Inches(4.0)
    accent_x = (SLIDE_W - accent_w) // 2
    _add_rect(slide, accent_x, Inches(4.4), accent_w, Inches(0.06), GOLD)

    # Course + professor block, centered (larger per round-3 feedback)
    _add_text(slide, MARGIN, Inches(4.8), RULE_W, Inches(0.55),
              "Management 405  ·  EMBA",
              size=26, bold=True, color=GRAY, font="Calibri",
              align=PP_ALIGN.CENTER)
    _add_text(slide, MARGIN, Inches(5.5), RULE_W, Inches(0.5),
              "Prof. Nico Voigtlaender  ·  UCLA Anderson",
              size=22, color=GRAY, font="Calibri",
              align=PP_ALIGN.CENTER)

    # Footer rhythm – keeps the cover in the same visual family
    _add_rect(slide, 0, Inches(7.15), SLIDE_W, Inches(0.02), RULE)
    _add_rect(slide, MARGIN, Inches(7.135), GOLD_W, Inches(0.05), GOLD)


# --------------------------------------------------------------------------
# Layout 2 – Section header (you-are-here agenda)
# --------------------------------------------------------------------------

# Module 3 agenda data. For each divider slide, one entry is "current".
MODULE_AGENDA = [
    {
        "title": "Part 1: Production – Picking the Right Inputs",
        "subs": [
            "Short Run – Hiring Decisions",
            "Long Run – Optimal Input Mix",
        ],
    },
    {
        "title": "Part 2: Costs – Producing at the Lowest Price",
        "subs": [
            "Cost Concepts",
            "Scale & Scope",
        ],
    },
]


def slide_section_header(prs, current_part_idx=0):
    """Render a section-divider slide ("Agenda" view).

    *current_part_idx* selects which agenda entry is "current" (navy).
    All other parts render in faded grey below. Subsections are real
    PPT auto-numbered list items, restarted at 1 within each Part.
    """
    slide = _blank_slide(prs)
    _draw_top_bar(slide, "Module 3 · Section Divider")
    _draw_action_title(slide, "Agenda")

    y = Inches(1.85)
    part_title_h = Inches(0.6)
    sub_list_h = Inches(1.35)        # comfortable space for 2 subs at 24pt
    block_gap = Inches(0.15)

    for idx, part in enumerate(MODULE_AGENDA):
        is_current = (idx == current_part_idx)
        color = NAVY if is_current else FADED

        # Part title
        _add_text(slide, MARGIN, y, RULE_W, part_title_h,
                  part["title"], size=30, bold=True, color=color,
                  font="Calibri")
        y += part_title_h

        # Subsections as a real PPT auto-numbered list (restarts at 1
        # for each Part because we use one text frame per Part)
        _add_bulleted_list(
            slide,
            left=MARGIN + Inches(0.4),
            top=y,
            width=RULE_W - Inches(0.4),
            height=sub_list_h,
            items=part["subs"],
            size=24, color=color, bullet_color=color,
            line_spacing_pts=10,
            autonum_scheme='alphaLcPeriod',
        )
        y += sub_list_h + block_gap

    _draw_footer(slide, FOOTER_TEXT, 6)


# --------------------------------------------------------------------------
# Layout 3 – Content bulleted (native PPT bullets, navy body + bullets)
# --------------------------------------------------------------------------

def slide_content_bulleted(prs):
    slide = _blank_slide(prs)
    _draw_top_bar(slide, "Module 3 · Part 1 · Production")
    _draw_action_title(slide, "Hire when MRPL > wage; stop when MRPL = wage")

    bullets = [
        "Add a worker → output rises by MPL",
        "Revenue from that worker = MRPL = MPL × price",
        "Hire if MRPL > wage; the next dollar earns more than it costs",
        "Stop when MRPL = wage – the marginal hire just breaks even",
    ]

    _add_bulleted_list(
        slide,
        left=MARGIN,
        top=Inches(1.85),
        width=RULE_W,
        height=Inches(5.0),
        items=bullets,
        size=24, color=NAVY, bullet_color=NAVY,
        line_spacing_pts=18,
    )

    _draw_footer(slide, FOOTER_TEXT, 21)


# --------------------------------------------------------------------------
# Layout 4 – Content two-column (the original sample, preserved)
# --------------------------------------------------------------------------

def slide_content_two_column(prs):
    slide = _blank_slide(prs)
    _draw_top_bar(slide, "Module 3 · Part 2 · Costs")
    _draw_action_title(slide, "Sunk costs should never drive decisions")

    rows = [
        ("Fixed costs", "don't depend on quantity produced (Q)"),
        ("Sunk costs", "a fixed cost that cannot be recovered"),
        ("→ Decision rule", "ignore sunk costs when choosing what to do next"),
        ("Variable costs", "depend on volume produced (Q)"),
    ]

    top = Inches(1.95)
    row_h = Inches(0.95)
    concept_w = Inches(3.2)
    col_gap = Inches(0.2)
    def_x = MARGIN + concept_w + col_gap
    def_w = RULE_W - concept_w - col_gap
    for i, (concept, defn) in enumerate(rows):
        y = top + row_h * i
        _add_text(slide, MARGIN, y, concept_w, row_h,
                  concept, size=18, bold=True, color=NAVY)
        _add_text(slide, def_x, y, def_w, row_h,
                  defn, size=18, color=GRAY)

    _draw_footer(slide, FOOTER_TEXT, 42)


# --------------------------------------------------------------------------
# Layout 5 – Poll slide (auto-num A./B./C./D. bullets in navy, no QR)
# --------------------------------------------------------------------------

def slide_poll(prs):
    slide = _blank_slide(prs)
    _draw_top_bar(slide, "Module 3 · Part 1 · Production")
    _draw_action_title(slide, "Tesla's MRPL at 6,000 employees – what is it?")
    _draw_poll_pill(slide)

    options = [
        "$25,000",
        "$50,000",
        "$75,000",
        "$100,000",
    ]
    _add_bulleted_list(
        slide,
        left=MARGIN,
        top=Inches(2.0),
        width=RULE_W,
        height=Inches(4.0),
        items=options,
        size=26, color=NAVY, bullet_color=NAVY,
        line_spacing_pts=20,
        autonum_scheme='alphaUcPeriod',
    )

    _add_text(slide, MARGIN, Inches(6.5), RULE_W, Inches(0.4),
              "Respond at PollEv.com/nvoigtlaender",
              size=14, italic=True, color=GRAY, align=PP_ALIGN.RIGHT)

    _draw_footer(slide, FOOTER_TEXT, 24)


# --------------------------------------------------------------------------
# Layout 6 – Closing synthesis
# --------------------------------------------------------------------------

def slide_closing_synthesis(prs):
    slide = _blank_slide(prs)
    _draw_top_bar(slide, "Module 3 · Synthesis & Bridge to Module 4")
    _draw_action_title(
        slide,
        "Production and costs set the stage for pricing and profit",
    )

    _add_text(slide, MARGIN, Inches(1.85), RULE_W, Inches(0.45),
              "MODULE 3 RECAP",
              size=14, bold=True, color=NAVY, font="Calibri")

    recap = [
        ("Production",
         "Hire until MRPL = wage; choose inputs by bang-for-the-buck"),
        ("Costs",
         "Ignore sunk; opportunity cost is real; watch MC vs. AC"),
        ("Scale",
         "Bigger is often cheaper – until diseconomies set in"),
    ]
    top = Inches(2.35)
    row_h = Inches(0.55)
    concept_w = Inches(2.2)
    col_gap = Inches(0.2)
    def_x = MARGIN + concept_w + col_gap
    def_w = RULE_W - concept_w - col_gap
    for i, (concept, defn) in enumerate(recap):
        y = top + row_h * i
        _add_text(slide, MARGIN, y, concept_w, row_h,
                  concept, size=18, bold=True, color=NAVY)
        _add_text(slide, def_x, y, def_w, row_h,
                  defn, size=18, color=GRAY)

    split_y = Inches(4.3)
    _add_rect(slide, MARGIN, split_y, RULE_W, Inches(0.04), GOLD)

    _add_text(slide, MARGIN, Inches(4.55), RULE_W, Inches(0.45),
              "COMING UP – MODULE 4: PRICING & PROFIT",
              size=14, bold=True, color=GOLD, font="Calibri")

    preview = [
        ("Combine",
         "demand (M2) with cost (M3) to find profit-maximizing price"),
        ("Decide",
         "how much to produce when each unit costs and earns something different"),
        ("Predict",
         "how price and quantity respond to cost shocks and demand shifts"),
    ]
    top2 = Inches(5.05)
    for i, (concept, defn) in enumerate(preview):
        y = top2 + row_h * i
        _add_text(slide, MARGIN, y, concept_w, row_h,
                  concept, size=18, bold=True, color=NAVY)
        _add_text(slide, def_x, y, def_w, row_h,
                  defn, size=18, color=GRAY)

    _draw_footer(slide, FOOTER_TEXT, 76)


# --------------------------------------------------------------------------
# Build orchestration
# --------------------------------------------------------------------------

def build_template():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_section_header(prs, current_part_idx=0)
    slide_content_bulleted(prs)
    slide_content_two_column(prs)
    slide_poll(prs)
    slide_closing_synthesis(prs)

    prs.save(OUT_DIR / "405 Slides Template.pptx")


if __name__ == "__main__":
    build_template()
    print(f"Wrote {OUT_DIR / '405 Slides Template.pptx'}")
