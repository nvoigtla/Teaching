"""Extract slide titles, text, and notes from Module 3.pptx for review."""
from pptx import Presentation
from pathlib import Path

deck = Path(__file__).parent / "Module 3.pptx"
prs = Presentation(deck)

out_path = Path(__file__).parent / "_outline_dump.txt"
lines = []
lines.append(f"Total slides: {len(prs.slides)}")
lines.append("=" * 80)

for i, slide in enumerate(prs.slides, 1):
    title = ""
    body_chunks = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = "\n".join(p.text for p in shape.text_frame.paragraphs).strip()
        if not text:
            continue
        if shape == slide.shapes.title:
            title = text
        else:
            body_chunks.append(text)

    notes = ""
    if slide.has_notes_slide:
        notes = slide.notes_slide.notes_text_frame.text.strip()

    layout = slide.slide_layout.name if slide.slide_layout else "?"

    lines.append("")
    lines.append(f"--- Slide {i} | Layout: {layout} ---")
    lines.append(f"TITLE: {title or '(no title)'}")
    if body_chunks:
        lines.append("BODY:")
        for chunk in body_chunks:
            for line in chunk.split("\n"):
                if line.strip():
                    lines.append(f"  - {line.strip()}")
    if notes:
        lines.append("NOTES:")
        notes_short = notes if len(notes) < 600 else notes[:600] + " [...]"
        for line in notes_short.split("\n"):
            if line.strip():
                lines.append(f"    {line.strip()}")

out_path.write_text("\n".join(lines), encoding="utf-8")
print(f"Wrote {out_path}")
