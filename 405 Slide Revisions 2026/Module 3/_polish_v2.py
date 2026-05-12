"""
Polish v2 for Module 3_NEW_draft.pptx:

  1. Force a WHITE background on the title slide (currently inherits a
     blue background from the master).
  2. Swap the title-slide image on the Apple Car slide (47) from Apple
     Park to a Vanarama Apple Car concept render, with adjusted
     dimensions to preserve aspect ratio and a new attribution caption.
  3. Add MBA-friendly speaker notes to every slide (76 total). For
     slides that already have a notesSlide, replace the body text. For
     slides that don't (the 7 new template-styled slides + slides with
     dropped rels), create a fresh notesSlide and wire it up.

Direct zip + lxml surgery throughout – python-pptx never sees the
working deck on the save path.
"""

import shutil
import sys
import zipfile
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches

HERE = Path(__file__).parent
sys.path.insert(0, str(HERE))
from _build_template_samples import (  # noqa: E402
    NAVY, GRAY, RULE, GOLD, WHITE,
    MARGIN, RULE_W, GOLD_W, SLIDE_W, SLIDE_H, FOOTER_TEXT,
    _add_text, _add_rect,
    _draw_top_bar, _draw_action_title, _draw_footer,
    _add_bulleted_list,
)
from _speaker_notes import NOTES  # noqa: E402

TARGET = HERE / "Module 3_NEW_draft.pptx"
APPLE_CAR_IMAGE = HERE / "_apple_car_concept.jpg"
SIDECAR_APPLE = HERE / "_sidecar_apple_v2.pptx"


NS = {
    'p':   "http://schemas.openxmlformats.org/presentationml/2006/main",
    'a':   "http://schemas.openxmlformats.org/drawingml/2006/main",
    'r':   "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    'rel': "http://schemas.openxmlformats.org/package/2006/relationships",
    'ct':  "http://schemas.openxmlformats.org/package/2006/content-types",
}
REL_SLIDE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide"
REL_LAYOUT = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout"
REL_IMAGE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"
REL_NOTES = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide"
REL_NOTES_MASTER = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesMaster"

CT_NOTES = "application/vnd.openxmlformats-officedocument.presentationml.notesSlide+xml"


def Q(prefix, tag):
    return f'{{{NS[prefix]}}}{tag}'


def parse(xml_bytes):
    return etree.fromstring(xml_bytes)


def serialize(tree):
    return etree.tostring(tree, xml_declaration=True, encoding='UTF-8', standalone=True)


def _resolve_relative(base_dir, rel_target):
    combined = base_dir + '/' + rel_target
    parts = []
    for seg in combined.replace('\\', '/').split('/'):
        if seg == '..':
            if parts:
                parts.pop()
        elif seg and seg != '.':
            parts.append(seg)
    return '/'.join(parts)


def get_slide_path_at_position(members, pos_1_based):
    pres = parse(members['ppt/presentation.xml'])
    pres_rels = parse(members['ppt/_rels/presentation.xml.rels'])
    sld_ids = pres.find(Q('p', 'sldIdLst')).findall(Q('p', 'sldId'))
    rId = sld_ids[pos_1_based - 1].get(Q('r', 'id'))
    for rel in pres_rels.findall(Q('rel', 'Relationship')):
        if rel.get('Id') == rId:
            return ('ppt/' + rel.get('Target')).replace('\\', '/')
    raise RuntimeError(f"no slide rel for rId={rId}")


def next_rId(rels_tree):
    used = {r.get('Id') for r in rels_tree.findall(Q('rel', 'Relationship'))}
    n = 1
    while f'rId{n}' in used:
        n += 1
    return f'rId{n}'


# ============================================================
# (1) White background on title slide
# ============================================================

def force_white_background_on_title(members):
    print("(1) force_white_background_on_title:")
    slide_path = get_slide_path_at_position(members, 1)
    print(f"    target slide = {slide_path}")
    slide = parse(members[slide_path])
    cSld = slide.find(Q('p', 'cSld'))

    # Remove existing <p:bg> if present
    for bg in cSld.findall(Q('p', 'bg')):
        cSld.remove(bg)

    # Build <p:bg><p:bgPr><a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill><a:effectLst/></p:bgPr></p:bg>
    bg = etree.Element(Q('p', 'bg'))
    bgPr = etree.SubElement(bg, Q('p', 'bgPr'))
    fill = etree.SubElement(bgPr, Q('a', 'solidFill'))
    srgb = etree.SubElement(fill, Q('a', 'srgbClr'))
    srgb.set('val', 'FFFFFF')
    etree.SubElement(bgPr, Q('a', 'effectLst'))

    # <p:bg> must be the FIRST child of <p:cSld>, before <p:spTree>
    cSld.insert(0, bg)

    members[slide_path] = serialize(slide)
    print(f"    white background set")


# ============================================================
# (2) Swap Apple Car image (Apple Park → Vanarama concept render)
# ============================================================

def _strip_slide(slide):
    spTree = slide.shapes._spTree
    for child in list(spTree):
        tag = child.tag.split('}', 1)[1]
        if tag in ('sp', 'grpSp', 'pic', 'graphicFrame', 'cxnSp', 'contentPart'):
            spTree.remove(child)
    slide._element.set('showMasterSp', '0')


def _new_blank_slide(prs):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _strip_slide(slide)
    return slide


def build_apple_sidecar_v2(image_path):
    prs = Presentation()
    prs.slide_width = Emu(int(SLIDE_W))
    prs.slide_height = Emu(int(SLIDE_H))
    slide = _new_blank_slide(prs)

    _draw_top_bar(slide, "Module 3 · Part 2 · Costs")
    _draw_action_title(slide,
        "Opportunity cost is a real cost: Apple's canceled Apple Car")

    # Vanarama concept render is 1500x808 → aspect 1.857
    # Display 5.6" wide × 3.01" tall to preserve aspect ratio
    img_w = Inches(5.6)
    img_h = Inches(3.01)
    img_left = Inches(7.4)
    img_top = Inches(2.2)
    slide.shapes.add_picture(str(image_path), img_left, img_top,
                             width=img_w, height=img_h)

    # Attribution
    _add_text(slide, img_left, Inches(5.30), img_w, Inches(0.45),
              "Vanarama Apple Car concept (fair use, © Vanarama)",
              size=10, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

    # Bullets (left ~52%)
    bullets = [
        "Apple killed Project Titan in 2024 after ~10 years and ~$10B spent",
        "Sunk costs ≠ a reason to keep going",
        "Real reason to stop: opportunity cost of capital + ~2,000 engineers",
        "Reallocated → AI / Apple Intelligence (the higher-MPL use)",
    ]
    _add_bulleted_list(
        slide,
        left=MARGIN, top=Inches(1.95),
        width=Inches(6.9), height=Inches(4.5),
        items=bullets,
        size=22, color=NAVY, bullet_color=NAVY,
        bullet_char="▪",
        bullet_size_pct=150,
        line_spacing_pts=18,
    )

    _draw_footer(slide, FOOTER_TEXT, 47)
    prs.save(SIDECAR_APPLE)


def swap_apple_car_image(members, image_path):
    print("(2) swap_apple_car_image:")
    if not image_path.exists():
        raise FileNotFoundError(f"image not found: {image_path}")

    build_apple_sidecar_v2(image_path)

    with zipfile.ZipFile(SIDECAR_APPLE, 'r') as zf:
        sidecar_slide_xml = zf.read('ppt/slides/slide1.xml')
        sidecar_slide_rels = zf.read('ppt/slides/_rels/slide1.xml.rels')
        media_files = [n for n in zf.namelist() if n.startswith('ppt/media/')]
        assert len(media_files) == 1
        sidecar_media_path = media_files[0]
        image_bytes = zf.read(sidecar_media_path)
        sidecar_ext = sidecar_media_path.rsplit('.', 1)[1]

    sidecar_rels_root = parse(sidecar_slide_rels)
    sidecar_image_rId = None
    for rel in sidecar_rels_root.findall(Q('rel', 'Relationship')):
        if rel.get('Type') == REL_IMAGE:
            sidecar_image_rId = rel.get('Id')
            break
    print(f"    sidecar image rId = {sidecar_image_rId}, ext = {sidecar_ext}")

    # Drop old Apple Park image
    old_path = 'ppt/media/_apple_park.jpg'
    if old_path in members:
        del members[old_path]
        print(f"    dropped old image at {old_path}")

    # Embed new image
    new_image_path = f'ppt/media/_apple_car.{sidecar_ext}'
    members[new_image_path] = image_bytes
    print(f"    embedded new image at {new_image_path} ({len(image_bytes):,} bytes)")

    # Ensure Content_Types has Default for the extension
    ct = parse(members['[Content_Types].xml'])
    if not any(d.get('Extension') == sidecar_ext for d in ct.findall(Q('ct', 'Default'))):
        default = etree.Element(Q('ct', 'Default'))
        default.set('Extension', sidecar_ext)
        default.set('ContentType', 'image/jpeg' if sidecar_ext in ('jpg', 'jpeg') else 'image/png')
        ct.insert(0, default)
        members['[Content_Types].xml'] = serialize(ct)

    # Locate Apple Car slide in destination
    slide_path = get_slide_path_at_position(members, 47)
    slide_dir, slide_file = slide_path.rsplit('/', 1)
    slide_rels_path = f'{slide_dir}/_rels/{slide_file}.rels'
    print(f"    target slide = {slide_path}")

    # Replace XML
    members[slide_path] = sidecar_slide_xml

    # Replace rels (layout + image)
    rels_root = etree.Element(Q('rel', 'Relationships'),
                              nsmap={None: NS['rel']})
    layout_rel = etree.SubElement(rels_root, Q('rel', 'Relationship'))
    layout_rel.set('Id', 'rId1')
    layout_rel.set('Type', REL_LAYOUT)
    layout_rel.set('Target', '../slideLayouts/slideLayout2.xml')
    image_rel = etree.SubElement(rels_root, Q('rel', 'Relationship'))
    image_rel.set('Id', sidecar_image_rId)
    image_rel.set('Type', REL_IMAGE)
    image_rel.set('Target', f'../media/_apple_car.{sidecar_ext}')
    members[slide_rels_path] = serialize(rels_root)
    print(f"    slide XML + rels updated")


# ============================================================
# (3) Speaker notes for all 76 slides
# ============================================================

def find_notes_master_target(members, from_dir='ppt/notesSlides'):
    """Return the Target string for the notesMaster relative to *from_dir*."""
    # notesMaster lives at ppt/notesMasters/notesMaster1.xml
    return '../notesMasters/notesMaster1.xml'


def next_notesSlide_filename(members):
    used = set()
    for name in members.keys():
        if name.startswith('ppt/notesSlides/notesSlide') and name.endswith('.xml'):
            tail = name[len('ppt/notesSlides/notesSlide'):-len('.xml')]
            try:
                used.add(int(tail))
            except ValueError:
                pass
    n = 1
    while n in used:
        n += 1
    return f'ppt/notesSlides/notesSlide{n}.xml'


def build_notes_slide_xml(note_text):
    """Construct a fresh notesSlide XML body containing *note_text*."""
    # Split note_text into paragraphs by double newlines, if any
    paragraphs = note_text.split('\n\n') if '\n\n' in note_text else [note_text]
    paragraphs_xml = '\n'.join(
        f'<a:p><a:r><a:rPr lang="en-US"/><a:t>{_xml_escape(p)}</a:t></a:r></a:p>'
        for p in paragraphs
    )
    return f'''<?xml version='1.0' encoding='UTF-8' standalone='yes'?>
<p:notes xmlns:a="{NS['a']}" xmlns:r="{NS['r']}" xmlns:p="{NS['p']}">
<p:cSld>
<p:spTree>
<p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>
<p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/><a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr>
<p:sp>
<p:nvSpPr><p:cNvPr id="2" name="Slide Image Placeholder 1"/><p:cNvSpPr><a:spLocks noGrp="1" noRot="1" noChangeAspect="1"/></p:cNvSpPr><p:nvPr><p:ph type="sldImg"/></p:nvPr></p:nvSpPr>
<p:spPr/>
</p:sp>
<p:sp>
<p:nvSpPr><p:cNvPr id="3" name="Notes Placeholder 2"/><p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr><p:ph type="body" idx="1"/></p:nvPr></p:nvSpPr>
<p:spPr/>
<p:txBody>
<a:bodyPr/>
<a:lstStyle/>
{paragraphs_xml}
</p:txBody>
</p:sp>
<p:sp>
<p:nvSpPr><p:cNvPr id="4" name="Slide Number Placeholder 3"/><p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr><p:ph type="sldNum" sz="quarter" idx="10"/></p:nvPr></p:nvSpPr>
<p:spPr/>
<p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:fld id="{{8D36AC2A-E969-3C46-9ED8-B46D40C5E669}}" type="slidenum"><a:rPr lang="en-US"/><a:t>1</a:t></a:fld><a:endParaRPr lang="en-US"/></a:p></p:txBody>
</p:sp>
</p:spTree>
</p:cSld>
<p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr>
</p:notes>'''


def _xml_escape(s):
    return (s.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;'))


def replace_notes_body(notes_xml_bytes, note_text):
    """In an existing notesSlide XML, replace the body placeholder's text."""
    root = parse(notes_xml_bytes)
    spTree = root.find(Q('p', 'cSld') + '/' + Q('p', 'spTree'))
    body_sp = None
    for sp in spTree.findall(Q('p', 'sp')):
        nvPr = sp.find(Q('p', 'nvSpPr') + '/' + Q('p', 'nvPr'))
        if nvPr is not None:
            ph = nvPr.find(Q('p', 'ph'))
            if ph is not None and ph.get('type') == 'body':
                body_sp = sp
                break
    if body_sp is None:
        return None
    txBody = body_sp.find(Q('p', 'txBody'))
    if txBody is None:
        # Add a fresh txBody
        txBody = etree.SubElement(body_sp, Q('p', 'txBody'))
        etree.SubElement(txBody, Q('a', 'bodyPr'))
        etree.SubElement(txBody, Q('a', 'lstStyle'))
    # Remove all existing <a:p> children
    for p_el in txBody.findall(Q('a', 'p')):
        txBody.remove(p_el)
    # Add a new <a:p> with the note text
    p = etree.SubElement(txBody, Q('a', 'p'))
    r = etree.SubElement(p, Q('a', 'r'))
    rPr = etree.SubElement(r, Q('a', 'rPr'))
    rPr.set('lang', 'en-US')
    t = etree.SubElement(r, Q('a', 't'))
    t.text = note_text
    return serialize(root)


def add_notes_to_slide(members, slide_position, note_text):
    """Add or replace speaker notes for the slide at *slide_position*."""
    slide_path = get_slide_path_at_position(members, slide_position)
    slide_dir, slide_file = slide_path.rsplit('/', 1)
    slide_rels_path = f'{slide_dir}/_rels/{slide_file}.rels'

    # 1) Does this slide already have a notesSlide rel?
    slide_rels = parse(members[slide_rels_path])
    notes_rel = None
    for rel in slide_rels.findall(Q('rel', 'Relationship')):
        if rel.get('Type') == REL_NOTES:
            notes_rel = rel
            break

    if notes_rel is not None:
        # Update existing notesSlide
        notes_target = notes_rel.get('Target')
        notes_path = _resolve_relative(slide_dir, notes_target)
        if notes_path not in members:
            print(f"    WARN: rel points to missing notes part {notes_path}")
            return
        new_xml = replace_notes_body(members[notes_path], note_text)
        if new_xml is None:
            print(f"    WARN: could not find body placeholder in {notes_path}")
            return
        members[notes_path] = new_xml
        return 'updated'
    else:
        # Create a new notesSlide
        new_notes_path = next_notesSlide_filename(members)
        new_notes_filename = new_notes_path.rsplit('/', 1)[1]
        new_notes_rels_path = f'ppt/notesSlides/_rels/{new_notes_filename}.rels'

        members[new_notes_path] = build_notes_slide_xml(note_text).encode('utf-8')

        # Build notesSlide rels: notesMaster + back to slide
        rels_root = etree.Element(Q('rel', 'Relationships'),
                                  nsmap={None: NS['rel']})
        master_rel = etree.SubElement(rels_root, Q('rel', 'Relationship'))
        master_rel.set('Id', 'rId1')
        master_rel.set('Type', REL_NOTES_MASTER)
        master_rel.set('Target', '../notesMasters/notesMaster1.xml')
        back_rel = etree.SubElement(rels_root, Q('rel', 'Relationship'))
        back_rel.set('Id', 'rId2')
        back_rel.set('Type', REL_SLIDE)
        back_rel.set('Target', f'../slides/{slide_file}')
        members[new_notes_rels_path] = serialize(rels_root)

        # Add Content_Types Override
        ct = parse(members['[Content_Types].xml'])
        ov = etree.SubElement(ct, Q('ct', 'Override'))
        ov.set('PartName', '/' + new_notes_path)
        ov.set('ContentType', CT_NOTES)
        members['[Content_Types].xml'] = serialize(ct)

        # Add rel from slide -> notesSlide
        new_rel_id = next_rId(slide_rels)
        new_rel = etree.SubElement(slide_rels, Q('rel', 'Relationship'))
        new_rel.set('Id', new_rel_id)
        new_rel.set('Type', REL_NOTES)
        new_rel.set('Target', f'../notesSlides/{new_notes_filename}')
        members[slide_rels_path] = serialize(slide_rels)

        return 'created'


def add_all_notes(members):
    print("(3) add_all_notes:")
    created = updated = 0
    for pos in sorted(NOTES.keys()):
        result = add_notes_to_slide(members, pos, NOTES[pos])
        if result == 'created':
            created += 1
        elif result == 'updated':
            updated += 1
    print(f"    notesSlides updated: {updated}, created: {created}")


# ============================================================
# Main
# ============================================================

def main():
    print(f"Reading {TARGET.name}")
    with zipfile.ZipFile(TARGET, 'r') as zf:
        members = {n: zf.read(n) for n in zf.namelist()}

    force_white_background_on_title(members)
    swap_apple_car_image(members, APPLE_CAR_IMAGE)
    add_all_notes(members)

    tmp = TARGET.with_suffix('.pptx.tmp')
    with zipfile.ZipFile(tmp, 'w', zipfile.ZIP_DEFLATED) as zf:
        for n, d in members.items():
            zf.writestr(n, d)
    tmp.replace(TARGET)
    print(f"\nWrote {TARGET.name}")


if __name__ == "__main__":
    main()
