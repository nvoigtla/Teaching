"""
Rebuild Module 3.pptx via direct zip + lxml surgery.

This script does NOT save the original deck through python-pptx — that
path silently drops "NULL" image rels which several slides depend on,
causing PowerPoint to repair the file to empty.

Strategy:
  1. Read Module 3.pptx (original) into memory as a dict {filename:bytes}.
  2. Build the 5 new template-styled slides in a CLEAN sidecar deck via
     python-pptx (round-trip on a fresh deck is safe – no NULL rels) and
     extract their slide XML + rels.
  3. Mutate the in-memory zip:
        - drop 10 cut slides (XML, rels, [Content_Types] Override,
          presentation.xml.rels Relationship, sldIdLst sldId, plus
          orphan notesSlide if any),
        - inject the 5 new slides at the end (assigning fresh slideN.xml
          filenames, rIds, and Override / Relationship entries),
        - reorder presentation.xml's sldIdLst per TARGET_ORDER,
        - overlay chrome (navy top bar + footer rhythm + page #) on
          every retained original slide,
        - overlay revised takeaway-titles on ~22 slides,
        - replace Tesla → Apple Car content on new slide 47.
  4. Write the result as Module 3_NEW_draft.pptx.

If anything fails, the original Module 3.pptx is untouched.
"""

import io
import shutil
import sys
import zipfile
from copy import deepcopy
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Emu, Inches

HERE = Path(__file__).parent
sys.path.insert(0, str(HERE))
from _build_template_samples import (  # noqa: E402
    NAVY, GRAY, FADED, RULE, GOLD, WHITE,
    MARGIN, RULE_W, GOLD_W, SLIDE_W, SLIDE_H, FOOTER_TEXT,
    MODULE_AGENDA,
    _add_text, _add_rect,
    _draw_top_bar, _draw_action_title, _draw_footer, _draw_poll_pill,
    _add_bulleted_list,
)

ORIGINAL = HERE / "Module 3.pptx"
BACKUP = HERE / "Module 3_backup_2026-05-11_pre-rebuild.pptx"
SIDECAR = HERE / "_new_slides_sidecar.pptx"
OUTPUT = HERE / "Module 3_NEW_draft.pptx"


# ============================================================
# OOXML namespaces
# ============================================================

NS = {
    'p':   "http://schemas.openxmlformats.org/presentationml/2006/main",
    'a':   "http://schemas.openxmlformats.org/drawingml/2006/main",
    'r':   "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    'rel': "http://schemas.openxmlformats.org/package/2006/relationships",
    'ct':  "http://schemas.openxmlformats.org/package/2006/content-types",
    'mc':  "http://schemas.openxmlformats.org/markup-compatibility/2006",
    'p14': "http://schemas.microsoft.com/office/powerpoint/2010/main",
}

CT_SLIDE = "application/vnd.openxmlformats-officedocument.presentationml.slide+xml"
CT_NOTES = "application/vnd.openxmlformats-officedocument.presentationml.notesSlide+xml"
REL_SLIDE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide"
REL_LAYOUT = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout"
REL_NOTES = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide"


def Q(prefix, tag):
    return f'{{{NS[prefix]}}}{tag}'


# ============================================================
# Target ordering & metadata
# ============================================================

TARGET_ORDER = [
    ("keep", 1),  ("keep", 2),  ("keep", 4),  ("keep", 5),  ("keep", 8),  ("keep", 6),
    ("new",  "section_part1"),
    ("keep", 10), ("keep", 11), ("keep", 12), ("keep", 14), ("keep", 16),
    ("keep", 18), ("keep", 19), ("keep", 20), ("keep", 21), ("keep", 22),
    ("keep", 23), ("keep", 24), ("keep", 25), ("keep", 26), ("keep", 27),
    ("keep", 29), ("keep", 30), ("keep", 31), ("keep", 32), ("keep", 33),
    ("keep", 34), ("keep", 35),
    ("new",  "section_part1_2"),
    ("keep", 37), ("keep", 38), ("keep", 39), ("keep", 40), ("keep", 41),
    ("keep", 42), ("keep", 43), ("keep", 44), ("keep", 45), ("keep", 46),
    ("keep", 47),
    ("new",  "section_part2"),
    ("keep", 49), ("keep", 50), ("keep", 51), ("keep", 52), ("keep", 53),
    ("keep", 54), ("keep", 55), ("keep", 56), ("keep", 57), ("keep", 58),
    ("keep", 59), ("keep", 60), ("keep", 61), ("keep", 62), ("keep", 63),
    ("keep", 64), ("keep", 65), ("keep", 66), ("keep", 67), ("keep", 68),
    ("new",  "section_part2_2"),
    ("keep", 70), ("keep", 71), ("keep", 72), ("keep", 73), ("keep", 74),
    ("keep", 75), ("keep", 76), ("keep", 77), ("keep", 78), ("keep", 79),
    ("keep", 80), ("keep", 81),
    ("new",  "closing"),
]

REVISED_TITLES = {
    5:  "Every executive decision is a production-and-cost decision",
    9:  "In the short run, you're stuck with your capacity",
    13: "Hire more labor, get less per worker: diminishing MPL",
    21: "Hire when MRPL > wage; stop when MRPL = wage",
    22: "The optimal hiring rule: MRPL = w",
    23: "Caution: wages are not always constant",
    24: "Big employers bid their own wages up",
    28: "Solution: marginal cost of the 3rd designer = $3M",
    29: "Are real-world wages = MRPL?",
    33: 'The "bang for the buck" rule: equalize MP per dollar',
    40: "When prices change, the input mix shifts: robot tax & union wages",
    43: "Sunk costs should never drive decisions",
    45: "Why studios finish movies they know will flop: Waterworld",
    47: "Opportunity cost is a real cost: Apple's canceled Apple Car",
    50: "Marginal cost ≠ average cost: Burn60 workout packages",
    53: "Marginal cost in finance: the true rate on a bigger loan",
    64: "In the long run, you can pick a better-sized plant",
    65: "Long-run AC = the best you can do once capacity is flexible",
    66: "Bigger usually means cheaper – economies of scale",
    69: "But bigger isn't always cheaper – diseconomies of scale",
    70: "Sharing capabilities across products: economies of scope",
    71: "Amazon: scale, scope, or both?",
}

# Slides where the original is a full-bleed picture (PollEv exports);
# we skip the title-overlay so we don't cover the question + QR.
POLL_NEW_INDICES = {19, 27, 38, 51, 58, 74, 75}

CUT_INDICES = [3, 7, 9, 13, 15, 17, 28, 36, 48, 69]   # 1-based originals


# ============================================================
# Build sidecar deck with the 5 new slides via python-pptx.
# A fresh deck has no NULL rels, so python-pptx round-trip is safe.
# ============================================================

def _strip_slide(slide):
    spTree = slide.shapes._spTree
    for child in list(spTree):
        tag = child.tag.split('}', 1)[1]
        if tag in ('sp', 'grpSp', 'pic', 'graphicFrame', 'cxnSp', 'contentPart'):
            spTree.remove(child)
    slide._element.set('showMasterSp', '0')


def _new_blank_slide(prs):
    layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(layout)
    _strip_slide(slide)
    return slide


def _build_section_header(prs, current_part_idx, page_no):
    slide = _new_blank_slide(prs)
    _draw_top_bar(slide, "Module 3 · Section Divider")
    _draw_action_title(slide, "Agenda")
    y = Inches(1.85)
    for idx, part in enumerate(MODULE_AGENDA):
        is_current = (idx == current_part_idx)
        color = NAVY if is_current else FADED
        _add_text(slide, MARGIN, y, RULE_W, Inches(0.6),
                  part["title"], size=30, bold=True, color=color, font="Calibri")
        y += Inches(0.6)
        _add_bulleted_list(
            slide, left=MARGIN + Inches(0.4), top=y,
            width=RULE_W - Inches(0.4), height=Inches(1.35),
            items=part["subs"],
            size=24, color=color, bullet_color=color,
            line_spacing_pts=10, autonum_scheme='alphaLcPeriod',
        )
        y += Inches(1.35) + Inches(0.15)
    _draw_footer(slide, FOOTER_TEXT, page_no)


def _build_closing(prs):
    slide = _new_blank_slide(prs)
    _draw_top_bar(slide, "Module 3 · Synthesis & Bridge to Module 4")
    _draw_action_title(slide,
        "Production and costs set the stage for pricing and profit")
    _add_text(slide, MARGIN, Inches(1.85), RULE_W, Inches(0.45),
              "MODULE 3 RECAP", size=14, bold=True, color=NAVY, font="Calibri")
    recap = [
        ("Production", "Hire until MRPL = wage; choose inputs by bang-for-the-buck"),
        ("Costs",      "Ignore sunk; opportunity cost is real; watch MC vs. AC"),
        ("Scale",      "Bigger is often cheaper – until diseconomies set in"),
    ]
    _emit_two_col(slide, recap, Inches(2.35))
    _add_rect(slide, MARGIN, Inches(4.3), RULE_W, Inches(0.04), GOLD)
    _add_text(slide, MARGIN, Inches(4.55), RULE_W, Inches(0.45),
              "COMING UP – MODULE 4: PRICING & PROFIT",
              size=14, bold=True, color=GOLD, font="Calibri")
    preview = [
        ("Combine", "demand (M2) with cost (M3) to find profit-maximizing price"),
        ("Decide",  "how much to produce when each unit costs and earns differently"),
        ("Predict", "how price and quantity respond to cost shocks and demand shifts"),
    ]
    _emit_two_col(slide, preview, Inches(5.05))
    _draw_footer(slide, FOOTER_TEXT, 76)


def _emit_two_col(slide, rows, top):
    row_h = Inches(0.55)
    concept_w = Inches(2.2)
    col_gap = Inches(0.2)
    def_x = MARGIN + concept_w + col_gap
    def_w = RULE_W - concept_w - col_gap
    for i, (concept, defn) in enumerate(rows):
        y = top + row_h * i
        _add_text(slide, MARGIN, y, concept_w, row_h, concept,
                  size=18, bold=True, color=NAVY)
        _add_text(slide, def_x, y, def_w, row_h, defn,
                  size=18, color=GRAY)


def build_sidecar():
    prs = Presentation()
    prs.slide_width = Emu(int(SLIDE_W))
    prs.slide_height = Emu(int(SLIDE_H))
    # 1: section Part 1
    _build_section_header(prs, current_part_idx=0, page_no=7)
    # 2: section §1.2  (still inside Part 1 – Part 1 stays navy)
    _build_section_header(prs, current_part_idx=0, page_no=30)
    # 3: section Part 2
    _build_section_header(prs, current_part_idx=1, page_no=42)
    # 4: section §2.2
    _build_section_header(prs, current_part_idx=1, page_no=63)
    # 5: closing synthesis
    _build_closing(prs)
    prs.save(SIDECAR)


# ============================================================
# Zip + XML helpers
# ============================================================

def read_zip(path):
    members = {}
    with zipfile.ZipFile(path, 'r') as zf:
        for name in zf.namelist():
            members[name] = zf.read(name)
    return members


def write_zip(path, members):
    with zipfile.ZipFile(path, 'w', zipfile.ZIP_DEFLATED) as zf:
        for name, data in members.items():
            zf.writestr(name, data)


def parse(xml_bytes):
    return etree.fromstring(xml_bytes)


def serialize(tree, *, xml_declaration=True, standalone=True):
    return etree.tostring(
        tree, xml_declaration=xml_declaration,
        encoding='UTF-8', standalone=standalone,
    )


# ============================================================
# Helpers that operate on the in-memory zip dict
# ============================================================

def slide_filename_for_rId(members, rId):
    """Given an rId in presentation.xml.rels, return the slide XML
    package path (e.g. 'ppt/slides/slide5.xml')."""
    rels = parse(members['ppt/_rels/presentation.xml.rels'])
    for rel in rels.findall(Q('rel', 'Relationship')):
        if rel.get('Id') == rId and rel.get('Type') == REL_SLIDE:
            target = rel.get('Target')
            return ('ppt/' + target).replace('\\', '/')
    return None


def list_original_slide_rIds(members):
    """Return [(rId, slide_path)] in the order they appear in sldIdLst."""
    pres = parse(members['ppt/presentation.xml'])
    rels = parse(members['ppt/_rels/presentation.xml.rels'])
    rel_target = {r.get('Id'): r.get('Target')
                  for r in rels.findall(Q('rel', 'Relationship'))
                  if r.get('Type') == REL_SLIDE}
    out = []
    for sld in pres.find(Q('p', 'sldIdLst')).findall(Q('p', 'sldId')):
        rId = sld.get(Q('r', 'id'))
        target = rel_target.get(rId)
        if target:
            out.append((rId, ('ppt/' + target).replace('\\', '/')))
    return out


def next_slide_filename(members, taken=None):
    """Return the next available ppt/slides/slideN.xml filename."""
    taken = taken or set()
    used = set()
    for name in list(members.keys()) + list(taken):
        if name.startswith('ppt/slides/slide') and name.endswith('.xml'):
            tail = name[len('ppt/slides/slide'):-len('.xml')]
            try:
                used.add(int(tail))
            except ValueError:
                pass
    n = 1
    while n in used:
        n += 1
    return f'ppt/slides/slide{n}.xml'


def next_rId(rels_tree):
    used = {r.get('Id') for r in rels_tree.findall(Q('rel', 'Relationship'))}
    n = 1
    while f'rId{n}' in used:
        n += 1
    return f'rId{n}'


def get_slide_notes_rId(slide_rels_tree):
    for rel in slide_rels_tree.findall(Q('rel', 'Relationship')):
        if rel.get('Type') == REL_NOTES:
            return rel.get('Target'), rel.get('Id')
    return None, None


# ============================================================
# Cut a slide and all its dependents (notesSlide if any)
# ============================================================

def cut_slide(members, slide_path):
    """Drop *slide_path* from the package: slide xml, slide rels,
    its notesSlide (if any), Content_Types Override, and the
    presentation-level Relationship + sldIdLst entry."""
    # 1) Find the slide's rels file, extract the notesSlide target
    slide_dir, slide_file = slide_path.rsplit('/', 1)
    slide_rels_path = f'{slide_dir}/_rels/{slide_file}.rels'
    notes_path = None
    notes_rels_path = None
    if slide_rels_path in members:
        srels = parse(members[slide_rels_path])
        notes_target, _ = get_slide_notes_rId(srels)
        if notes_target:
            # Target like "../notesSlides/notesSlide5.xml"
            notes_path = _resolve_relative(slide_dir, notes_target)
            notes_dir, notes_file = notes_path.rsplit('/', 1)
            notes_rels_path = f'{notes_dir}/_rels/{notes_file}.rels'
    # 2) Find the presentation-level rId for this slide
    pres_rels = parse(members['ppt/_rels/presentation.xml.rels'])
    rId_to_drop = None
    for rel in list(pres_rels.findall(Q('rel', 'Relationship'))):
        if rel.get('Type') == REL_SLIDE:
            t = ('ppt/' + rel.get('Target')).replace('\\', '/')
            if t == slide_path:
                rId_to_drop = rel.get('Id')
                pres_rels.remove(rel)
                break
    members['ppt/_rels/presentation.xml.rels'] = serialize(pres_rels)

    # 3) Drop the sldId entry from presentation.xml
    pres = parse(members['ppt/presentation.xml'])
    sldIdLst = pres.find(Q('p', 'sldIdLst'))
    for sld in list(sldIdLst.findall(Q('p', 'sldId'))):
        if sld.get(Q('r', 'id')) == rId_to_drop:
            sldIdLst.remove(sld)
    members['ppt/presentation.xml'] = serialize(pres)

    # 4) Drop the slide file(s)
    for p in (slide_path, slide_rels_path):
        members.pop(p, None)

    # 5) Drop notesSlide files if they exist and aren't referenced
    #    by any other slide.  We trust that each notesSlide belongs
    #    to exactly one regular slide in the original deck.
    if notes_path:
        members.pop(notes_path, None)
    if notes_rels_path:
        members.pop(notes_rels_path, None)

    # 6) Update [Content_Types].xml – drop the Override for the slide
    #    (and the notesSlide if dropped).
    ct = parse(members['[Content_Types].xml'])
    parts_to_drop = {slide_path}
    if notes_path:
        parts_to_drop.add(notes_path)
    for ov in list(ct.findall(Q('ct', 'Override'))):
        pn = ov.get('PartName', '').lstrip('/')
        if pn in parts_to_drop:
            ct.remove(ov)
    members['[Content_Types].xml'] = serialize(ct)


def _resolve_relative(base_dir, rel_target):
    """Resolve a relationship Target (which is relative to the rels
    file's part directory) into a package-absolute path."""
    combined = base_dir + '/' + rel_target
    parts = []
    for seg in combined.replace('\\', '/').split('/'):
        if seg == '..':
            if parts:
                parts.pop()
        elif seg and seg != '.':
            parts.append(seg)
    return '/'.join(parts)


# ============================================================
# Inject a new slide into the package.
#
#   The slide XML and its rels XML come from the sidecar deck.
#   We rewrite the layout target so it points to slideLayout1 in the
#   destination (the "Title Only" layout in Module 3).
# ============================================================

def inject_slide(members, slide_xml_bytes, layout_target):
    new_slide_path = next_slide_filename(members)
    slide_dir, slide_file = new_slide_path.rsplit('/', 1)
    new_rels_path = f'{slide_dir}/_rels/{slide_file}.rels'

    # 1) Build the slide rels (just a layout reference)
    rels_root = etree.Element(Q('rel', 'Relationships'),
                              nsmap={None: NS['rel']})
    layout_rel = etree.SubElement(rels_root, Q('rel', 'Relationship'))
    layout_rel.set('Id', 'rId1')
    layout_rel.set('Type', REL_LAYOUT)
    layout_rel.set('Target', layout_target)
    members[new_rels_path] = serialize(rels_root)

    # 2) Add the slide XML
    members[new_slide_path] = slide_xml_bytes

    # 3) Add a presentation-level Relationship and a fresh rId
    pres_rels = parse(members['ppt/_rels/presentation.xml.rels'])
    rId = next_rId(pres_rels)
    new_rel = etree.SubElement(pres_rels, Q('rel', 'Relationship'))
    new_rel.set('Id', rId)
    new_rel.set('Type', REL_SLIDE)
    # Target relative to ppt/ – e.g. "slides/slide82.xml"
    new_rel.set('Target', new_slide_path[len('ppt/'):])
    members['ppt/_rels/presentation.xml.rels'] = serialize(pres_rels)

    # 4) Append a sldId to presentation.xml's sldIdLst (we'll reorder later)
    pres = parse(members['ppt/presentation.xml'])
    sldIdLst = pres.find(Q('p', 'sldIdLst'))
    existing_ids = [int(s.get('id'))
                    for s in sldIdLst.findall(Q('p', 'sldId'))]
    new_id = max(existing_ids) + 1 if existing_ids else 256
    new_sld = etree.SubElement(sldIdLst, Q('p', 'sldId'))
    new_sld.set('id', str(new_id))
    new_sld.set(Q('r', 'id'), rId)
    members['ppt/presentation.xml'] = serialize(pres)

    # 5) Register the new slide in [Content_Types].xml
    ct = parse(members['[Content_Types].xml'])
    ov = etree.SubElement(ct, Q('ct', 'Override'))
    ov.set('PartName', '/' + new_slide_path)
    ov.set('ContentType', CT_SLIDE)
    members['[Content_Types].xml'] = serialize(ct)

    return new_slide_path, rId


# ============================================================
# Reorder sldIdLst per TARGET_ORDER
# ============================================================

def reorder_sldIdLst(members, target_rId_sequence):
    pres = parse(members['ppt/presentation.xml'])
    sldIdLst = pres.find(Q('p', 'sldIdLst'))
    sld_by_rId = {s.get(Q('r', 'id')): s for s in sldIdLst.findall(Q('p', 'sldId'))}
    # Remove all
    for s in list(sldIdLst.findall(Q('p', 'sldId'))):
        sldIdLst.remove(s)
    # Re-append in target order
    for rId in target_rId_sequence:
        sld = sld_by_rId[rId]
        sldIdLst.append(sld)
    members['ppt/presentation.xml'] = serialize(pres)


# ============================================================
# Slide-level XML modifications – chrome overlay, revised titles,
# Apple Car content swap.
# ============================================================

def _make_rect_sp_xml(shp_id, x_emu, y_emu, w_emu, h_emu, hex_color):
    return f'''
<p:sp xmlns:p="{NS['p']}" xmlns:a="{NS['a']}">
  <p:nvSpPr>
    <p:cNvPr id="{shp_id}" name="Chrome{shp_id}"/>
    <p:cNvSpPr/><p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="{x_emu}" y="{y_emu}"/><a:ext cx="{w_emu}" cy="{h_emu}"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="{hex_color}"/></a:solidFill>
    <a:ln><a:noFill/></a:ln>
  </p:spPr>
  <p:txBody>
    <a:bodyPr/><a:lstStyle/><a:p/>
  </p:txBody>
</p:sp>'''


def _make_text_sp_xml(shp_id, x_emu, y_emu, w_emu, h_emu, text,
                      size_pt=18, bold=False, color_hex="000000",
                      align="l", anchor="t"):
    bold_attr = ' b="1"' if bold else ''
    return f'''
<p:sp xmlns:p="{NS['p']}" xmlns:a="{NS['a']}">
  <p:nvSpPr>
    <p:cNvPr id="{shp_id}" name="ChromeText{shp_id}"/>
    <p:cNvSpPr txBox="1"/><p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="{x_emu}" y="{y_emu}"/><a:ext cx="{w_emu}" cy="{h_emu}"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:noFill/>
  </p:spPr>
  <p:txBody>
    <a:bodyPr wrap="square" lIns="0" rIns="0" tIns="0" bIns="0" anchor="{anchor}"><a:spAutoFit/></a:bodyPr>
    <a:lstStyle/>
    <a:p>
      <a:pPr algn="{align}"/>
      <a:r>
        <a:rPr lang="en-US" sz="{int(size_pt * 100)}"{bold_attr}>
          <a:solidFill><a:srgbClr val="{color_hex}"/></a:solidFill>
          <a:latin typeface="Calibri"/>
        </a:rPr>
        <a:t>{_xml_escape(text)}</a:t>
      </a:r>
    </a:p>
  </p:txBody>
</p:sp>'''


def _xml_escape(s):
    return (s.replace('&', '&amp;')
             .replace('<', '&lt;')
             .replace('>', '&gt;'))


# Geometry constants in EMU
_EMU = 914400  # per inch
SLIDE_W_EMU = int(SLIDE_W)
MARGIN_EMU = int(MARGIN)
RULE_W_EMU = int(RULE_W)
GOLD_W_EMU = int(GOLD_W)
TOP_BAR_H = int(Inches(0.34))


def _next_shape_id(spTree):
    used = set()
    for cnv in spTree.iter(Q('p', 'cNvPr')):
        try:
            used.add(int(cnv.get('id')))
        except (TypeError, ValueError):
            pass
    n = 100  # start high to avoid clashes
    while n in used:
        n += 1
    return n


def apply_chrome_overlay(members, slide_path, page_num):
    """Add navy top bar + footer rhythm + page number to *slide_path*."""
    slide = parse(members[slide_path])
    spTree = slide.find(Q('p', 'cSld') + '/' + Q('p', 'spTree'))
    nid = _next_shape_id(spTree)

    # Navy top bar
    spTree.append(etree.fromstring(_make_rect_sp_xml(
        nid, 0, 0, SLIDE_W_EMU, TOP_BAR_H, "0B2B4E")))
    # Module tag on top bar
    spTree.append(etree.fromstring(_make_text_sp_xml(
        nid + 1, MARGIN_EMU, 0, int(Inches(12)), TOP_BAR_H,
        "MODULE 3  ·  PRODUCTION AND COSTS",
        size_pt=11, bold=True, color_hex="FFFFFF", anchor="ctr")))
    # Footer rule
    spTree.append(etree.fromstring(_make_rect_sp_xml(
        nid + 2, 0, int(Inches(7.15)), SLIDE_W_EMU, int(Inches(0.02)),
        "C8CDD3")))
    # Gold accent on footer
    spTree.append(etree.fromstring(_make_rect_sp_xml(
        nid + 3, MARGIN_EMU, int(Inches(7.135)), GOLD_W_EMU,
        int(Inches(0.05)), "E09F3E")))
    # Footer text
    spTree.append(etree.fromstring(_make_text_sp_xml(
        nid + 4, MARGIN_EMU, int(Inches(7.22)), int(Inches(11)),
        int(Inches(0.3)), FOOTER_TEXT,
        size_pt=10, color_hex="555B66")))
    # Page number
    spTree.append(etree.fromstring(_make_text_sp_xml(
        nid + 5, int(Inches(12.6)), int(Inches(7.22)),
        int(Inches(0.5)), int(Inches(0.3)), str(page_num),
        size_pt=10, color_hex="555B66", align="r")))

    members[slide_path] = serialize(slide)


def apply_revised_title(members, slide_path, new_title):
    """Overlay a white title band + new takeaway title between the top
    bar and the existing slide content."""
    slide = parse(members[slide_path])
    spTree = slide.find(Q('p', 'cSld') + '/' + Q('p', 'spTree'))
    nid = _next_shape_id(spTree)

    # White band to cover any existing title
    spTree.append(etree.fromstring(_make_rect_sp_xml(
        nid, 0, TOP_BAR_H, SLIDE_W_EMU, int(Inches(0.7)), "FFFFFF")))
    # New title text
    spTree.append(etree.fromstring(_make_text_sp_xml(
        nid + 1, MARGIN_EMU, int(Inches(0.4)),
        RULE_W_EMU, int(Inches(0.65)),
        new_title, size_pt=26, bold=True, color_hex="0B2B4E")))
    # Gold accent rail under title
    spTree.append(etree.fromstring(_make_rect_sp_xml(
        nid + 2, MARGIN_EMU, int(Inches(1.06)),
        RULE_W_EMU, int(Inches(0.02)), "C8CDD3")))
    spTree.append(etree.fromstring(_make_rect_sp_xml(
        nid + 3, MARGIN_EMU, int(Inches(1.045)),
        GOLD_W_EMU, int(Inches(0.05)), "E09F3E")))

    members[slide_path] = serialize(slide)


def apply_apple_car_content(members, slide_path):
    """Replace Tesla expansion-choice with Apple Car content.
    We blank the slide body (white rect across center) and overlay
    new takeaway bullets — non-destructive to chrome layers below."""
    slide = parse(members[slide_path])
    spTree = slide.find(Q('p', 'cSld') + '/' + Q('p', 'spTree'))
    nid = _next_shape_id(spTree)

    # Big white rectangle covering the body area
    spTree.append(etree.fromstring(_make_rect_sp_xml(
        nid, 0, int(Inches(1.3)), SLIDE_W_EMU,
        int(Inches(5.8)), "FFFFFF")))

    # Bullets as separate text boxes (no native PPT bullets here –
    # keeps the XML straightforward)
    bullets = [
        "Apple killed Project Titan in 2024 after ~10 years and ~$10B spent",
        "Sunk costs ≠ a reason to keep going",
        "Real reason to stop: opportunity cost of capital + ~2,000 engineers",
        "Reallocated → AI / Apple Intelligence (the higher-MPL use)",
    ]
    y0 = int(Inches(2.0))
    row_h = int(Inches(0.8))
    for i, b in enumerate(bullets):
        # gold square bullet
        spTree.append(etree.fromstring(_make_rect_sp_xml(
            nid + 1 + 2 * i,
            MARGIN_EMU, y0 + row_h * i + int(Inches(0.18)),
            int(Inches(0.16)), int(Inches(0.16)), "E09F3E")))
        # bullet text
        spTree.append(etree.fromstring(_make_text_sp_xml(
            nid + 2 + 2 * i,
            MARGIN_EMU + int(Inches(0.4)),
            y0 + row_h * i,
            RULE_W_EMU - int(Inches(0.4)), row_h,
            b, size_pt=22, color_hex="0B2B4E")))

    members[slide_path] = serialize(slide)


# ============================================================
# Main
# ============================================================

def main():
    print(f"Backup: {ORIGINAL.name} -> {BACKUP.name}")
    if not BACKUP.exists():
        shutil.copy(ORIGINAL, BACKUP)

    print(f"Building sidecar deck of new slides: {SIDECAR.name}")
    build_sidecar()

    print("Extracting new slide XMLs from sidecar")
    sidecar = read_zip(SIDECAR)
    # The sidecar has 5 slides: slide1..slide5
    new_slide_xmls = {}
    new_slide_xmls['section_part1']   = sidecar['ppt/slides/slide1.xml']
    new_slide_xmls['section_part1_2'] = sidecar['ppt/slides/slide2.xml']
    new_slide_xmls['section_part2']   = sidecar['ppt/slides/slide3.xml']
    new_slide_xmls['section_part2_2'] = sidecar['ppt/slides/slide4.xml']
    new_slide_xmls['closing']         = sidecar['ppt/slides/slide5.xml']

    print(f"Reading original: {ORIGINAL.name}")
    members = read_zip(ORIGINAL)
    original_slides = list_original_slide_rIds(members)
    print(f"  {len(original_slides)} slides in original deck")

    # Map 1-based original index -> rId (BEFORE any cuts)
    orig_idx_to_rId = {i + 1: rId for i, (rId, _) in enumerate(original_slides)}
    orig_idx_to_path = {i + 1: path for i, (_, path) in enumerate(original_slides)}

    # Cut slides
    for idx in CUT_INDICES:
        path = orig_idx_to_path[idx]
        print(f"  CUT slide #{idx}: {path}")
        cut_slide(members, path)

    # Insert new slides
    new_rIds = {}
    for name in ('section_part1', 'section_part1_2', 'section_part2',
                 'section_part2_2', 'closing'):
        path, rId = inject_slide(members, new_slide_xmls[name],
                                 layout_target='../slideLayouts/slideLayout2.xml')
        new_rIds[name] = rId
        print(f"  INSERT new {name!r} as {path} (rId={rId})")

    # Build target sequence of rIds
    target_seq = []
    for kind, ref in TARGET_ORDER:
        if kind == 'keep':
            target_seq.append(orig_idx_to_rId[ref])
        else:
            target_seq.append(new_rIds[ref])
    assert len(target_seq) == len(set(target_seq)), "duplicate rIds in target!"
    print(f"  reordering sldIdLst – target length {len(target_seq)}")
    reorder_sldIdLst(members, target_seq)

    # Apply chrome + revised titles + Apple Car content per the new
    # positions in TARGET_ORDER. We need to know each slide's package
    # path under the new ordering.
    # Build: new_pos (1-based) -> slide_path
    pres = parse(members['ppt/presentation.xml'])
    pres_rels = parse(members['ppt/_rels/presentation.xml.rels'])
    rId_to_target = {r.get('Id'): ('ppt/' + r.get('Target')).replace('\\', '/')
                     for r in pres_rels.findall(Q('rel', 'Relationship'))
                     if r.get('Type') == REL_SLIDE}
    sld_seq = pres.find(Q('p', 'sldIdLst')).findall(Q('p', 'sldId'))
    new_pos_to_path = {}
    new_pos_to_kind = {}
    for i, sld in enumerate(sld_seq, start=1):
        rId = sld.get(Q('r', 'id'))
        new_pos_to_path[i] = rId_to_target[rId]
        # kind: 'new' if this rId is one of the injected new slides
        new_pos_to_kind[i] = ('new' if rId in new_rIds.values() else 'keep')

    print(f"  applying chrome + titles + Apple Car content")
    n_chrome = 0
    n_titled = 0
    for new_pos in range(1, len(sld_seq) + 1):
        if new_pos_to_kind[new_pos] == 'new':
            # Section dividers and closing already carry chrome.
            continue
        path = new_pos_to_path[new_pos]
        apply_chrome_overlay(members, path, new_pos)
        n_chrome += 1
        if new_pos in REVISED_TITLES and new_pos not in POLL_NEW_INDICES:
            apply_revised_title(members, path, REVISED_TITLES[new_pos])
            n_titled += 1

    # Apple Car content swap on new slide 47
    apple_path = new_pos_to_path[47]
    apply_apple_car_content(members, apple_path)
    print(f"  Apple Car content applied on new slide 47 ({apple_path})")
    print(f"  total chrome overlays: {n_chrome}, revised titles: {n_titled}")

    # Write final zip
    if OUTPUT.exists():
        OUTPUT.unlink()
    write_zip(OUTPUT, members)
    print(f"\nWrote {OUTPUT.name}")


if __name__ == "__main__":
    main()
