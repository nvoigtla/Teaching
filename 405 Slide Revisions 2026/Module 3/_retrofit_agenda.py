# -*- coding: utf-8 -*-
"""Retrofit Module 3's agenda / divider slides to the numbered-circle
outline format (Teaching CLAUDE.md "Module-Outline / Agenda Slides"),
approved by Nico 2026-08-20.

Module 3's build script is FROZEN (deck = source of truth), so this is
in-place OOXML surgery: the replacement slides are generated here in a
temp deck (using Module 2's helper layer for the primitives — chrome,
bullets, footer, palette), then each agenda slide's <p:spTree> is
swapped for the generated one. Slide rels stay untouched (the new
shapes reference no rels); none of the agenda slides carries
animations (verified). Hidden slides elsewhere in the deck (displays
2–4) are not affected.

2026-08-26 (Nico), round 1 — brought up to the current outline
  standard:
  * non-current items are SHADED (#BFBFBF circle digit + item title);
    the gold circle fill stays gold on every item;
  * a shaded one-line row is nudged down by DIM_DROP so it centres in
    its reserved two-row box.
  Coverage pills were held back at that point (the deck was going to
  colleagues who do not tape videos); they are ON as of round 4 below.
  The maker is LOCAL here rather than a call into make_m2_outline,
  which expects Module 2's four-tuple outline rows.

2026-08-26 (Nico), round 2 — the agenda sequence itself:
  * the FIRST agenda slide (display 8) is now the descriptive overview
    of the whole module: every item lit, every description shown, no
    cream band (the standard's "descriptive overview");
  * a NEW slide is inserted right after it that highlights item 1 only,
    so the module opens "here is the plan → here is where we start";
  * the Part-2 divider (was display 41, now 42) highlights item 4 only
    — item 5 has its own agenda slide further down.
  Slide 8's old narration ("Let's start with Part 1, Production…")
  moved verbatim to the new item-1 slide; slide 8 got fresh
  whole-module overview narration.

2026-08-26 (Nico), round 3:
  * "Wage Searchers" joins the outline as item 3 (slides 26–31 taught
    it but no agenda item named it), so the module now has SIX items;
  * its agenda slide is inserted after the last short-run slide
    ("Optimal Hiring Rule in the Short Run: Numerical Solution");
  * TOP-BAR CONVENTION: every content slide's section tag now names
    the agenda item it belongs to — "Module 3 · Short Run: Hiring
    Decisions" — and agenda slides read "Module 3 · Agenda". Slides
    ahead of the first agenda slide (logistics, announcements, recap,
    course roadmap, big picture, concept map) keep their own tags, as
    does the summary closer. Written up in Teaching/CLAUDE.md.

2026-08-26 (Nico), round 4 — video conversion, step 2 of 2: every
  outline item carries its COVERAGE PILL at the right of the row (see
  COVERAGE_LABEL). Module 3 is taped end to end, so all six are gold
  "Video N" pills; the pill dims with its row on a section agenda. Step
  1, a title card per taped section, is _video_prep.py — run that first.

Rerunnable: the agenda slides are located by their on-slide text
("Outline of Module 3"), not by hard-coded display numbers; the insert
only happens when the deck still has one agenda slide fewer than
SPECS; and the retag pass is idempotent. Every run regenerates all
agenda spTrees from scratch.
"""
import os
import re
import sys
import zipfile
import shutil
from pathlib import Path

from lxml import etree as ET

sys.stdout.reconfigure(encoding='utf-8')
HERE = Path(__file__).parent
DECK = HERE / "Module 3 - Revised.pptx"
M2_DIR = HERE.parent / "Module 2"
sys.path.insert(0, str(M2_DIR))

import _build_Module2InClass as M2  # noqa: E402  (helper layer)

from pptx.enum.shapes import MSO_SHAPE  # noqa: E402
from pptx.enum.text import PP_ALIGN  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402

A = "http://schemas.openxmlformats.org/drawingml/2006/main"
P = "http://schemas.openxmlformats.org/presentationml/2006/main"
R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
PKG_R = "http://schemas.openxmlformats.org/package/2006/relationships"
CT = "http://schemas.openxmlformats.org/package/2006/content-types"


def q(ns, t):
    return "{%s}%s" % (ns, t)


# --- the outline lives in _m3_outline.py (shared with _video_prep.py) ---
from _m3_outline import (M3_OUTLINE, COVERAGE_LABEL,  # noqa: E402
                         IN_CLASS_ITEMS)

# 2026-08-26 (Nico, hand-set on the slide-11 agenda): the coverage pill
# sits vertically CENTRED in its row — 0.16" below the top of the
# reserved two-row box, which centres the 0.36" pill in the 0.715" box.
# A section agenda shades most rows down to a single centred line, so a
# top-aligned pill would float away from its item; the same drop is used
# on the descriptive overview too, so the pill column is identical on
# every agenda slide.
PILL_DROP = 146304                        # 0.16"

# Every description has to clear the pill by at least a five-letter
# word, or it reads as running into it (Nico, 2026-08-26, on the Cost
# Concepts line reaching under the Video 6 pill). The text box starts at
# x 2.05" and the pill's left edge is at 11.30", so the check is done in
# the real font — shorten the description, never shrink the type.
DESC_X = 2.05
PILL_LEFT = 11.30
CLEARANCE_WORD = "costs"                  # a five-letter word at 22 pt


def _check_descriptions():
    from PIL import ImageFont
    font = ImageFont.truetype("C:/Windows/Fonts/calibri.ttf", 22 * 4)

    def width(s):
        return font.getlength(s) / 4 / 72
    limit = PILL_LEFT - DESC_X - width(CLEARANCE_WORD)
    over = [(i + 1, width(d), d) for i, (_, d) in enumerate(M3_OUTLINE)
            if width(d) > limit]
    if over:
        for n, w, d in over:
            print("item %d's description runs %.2f\" (max %.2f\"): %s"
                  % (n, w, limit, d))
        raise SystemExit("shorten the description(s) above — they reach "
                         "under the coverage pill")
    print("descriptions clear the pills (max %.2f\", widest %.2f\")"
          % (limit, max(width(d) for _, d in M3_OUTLINE)))


_check_descriptions()

TAG = "Module 3 · Agenda"
TITLE = "Outline of Module 3"
FOOTER = "Management 405  ·  Module 3  ·  Production and Costs"

# The agenda sequence, in deck order.
# Item 0, the Introduction, has no agenda slide of its own: the module
# overview IS its agenda, so the section agendas run from item 1.
SPECS = [{"descriptions": True}] + [
    {"highlight": {i}} for i in range(1, len(M3_OUTLINE))]

# The one agenda slide this script may still have to CREATE: index into
# SPECS, the content slide it follows (matched on unique on-slide text),
# and its narration.
INSERT_SPEC = 4    # the Wage Searchers agenda
INSERT_AFTER = "Numerical Solution"
INSERT_NOTES = (
    "One more wrinkle before we leave the short run. Everything so far "
    "assumed the firm can hire as many people as it wants at the going "
    "wage. That is fine for a small employer. It is not fine for a firm "
    "that is large in its local labor market, or that is bidding for "
    "very scarce talent. There the next hire raises the wage for "
    "everyone, and the true cost of that hire is well above the wage on "
    "the offer letter. That is what a wage searcher faces, and it is "
    "what the next few slides work out."
)

# --- top-bar convention (round 3) -----------------------------------------
# Content slides carry the agenda item they belong to; agenda slides read
# TAG.  Slides ahead of the first agenda slide keep their own tags, and
# so does anything listed in TOPBAR_KEEP.
TOPBAR_PREFIX = "Module 3 · "
TOPBAR_KEEP = {"Module 3 · Summary"}
# front-matter tags that still carried the old three-level form
FRONT_TAG_FIX = {"Module 3 · Production · Big Picture":
                 "Module 3 · Big Picture"}


def make_m3_outline(prs, page_num, *, section_tag=TAG, title=TITLE,
                    descriptions=False, highlight_set=None):
    """One Module 3 agenda slide. Geometry is the Module 2 reference
    (0.42 title row + 0.38 description row + 0.11 gap = 0.91" pitch,
    the block centred between y 1.60" and 7.02"), so item positions are
    pixel-identical across all agenda slides. Every item RESERVES the
    description row; the description shows only for the current
    topic(s), or for all items when descriptions=True."""
    slide = M2._blank_slide(prs)
    M2._draw_top_bar_tc(slide, section_tag)
    M2._draw_action_title(slide, title)

    hi = set(highlight_set or ())
    if descriptions:
        hi = set(range(len(M3_OUTLINE)))

    # 2026-08-26: with the Introduction added the outline runs to SEVEN
    # items and the Module 2 metrics (0.42 + 0.38 + 0.11 = 0.91" pitch)
    # no longer fit between the title rule and the footer. The pitch is
    # therefore derived from the space available and capped at 0.91", so
    # a six-item deck keeps the Module 2 geometry to the EMU. Ported
    # verbatim from make_m1_outline, which solved this for Module 1's
    # seven items; the cap on the last row keeps its coverage pill clear
    # of the footer.
    n_items = len(M3_OUTLINE)
    top = Inches(1.42)
    bottom = Inches(7.02)
    gap = Inches(0.11) if n_items <= 6 else Inches(0.07)
    PITCH_MAX = Inches(0.91)             # the Module 2 pitch
    LAST_ROW_MAX = Inches(6.22)          # pill bottom then lands at 6.60"

    def _row_y(pitch, i):
        block = pitch * n_items - gap
        y0 = top + max(0, (bottom - top - block) // 2)
        return y0 + i * pitch

    pitch = PITCH_MAX
    while pitch > Inches(0.60) and (
            _row_y(pitch, n_items - 1) > LAST_ROW_MAX
            or pitch * n_items - gap > bottom - top):
        pitch -= 4572                     # 0.005" steps
    content = pitch - gap
    title_h = int(content * 0.525)        # 0.42 / 0.80 at the Module 2 pitch
    desc_h = content - title_h
    y = int(_row_y(pitch, 0))

    for i, (item, desc) in enumerate(M3_OUTLINE):
        # items not currently covered are shaded; a slide that lights
        # every item (descriptions=True) keeps them all navy
        lit = descriptions or i in hi
        ink = M2.NAVY if lit else M2.DIM
        if not descriptions and i in hi:
            band = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, int(Inches(0.90)),
                int(y - Inches(0.06)), int(Inches(12.15)),
                int(title_h + desc_h + Inches(0.10)))
            band.adjustments[0] = 0.35
            band.fill.solid()
            band.fill.fore_color.rgb = M2.CREAM
            band.line.color.rgb = M2.GOLD
            band.line.width = Pt(1.0)
            band.shadow.inherit = False
            M2._add_drop_shadow(band)
        circ = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, int(Inches(1.15)), int(y + Inches(0.02)),
            int(Inches(0.58)), int(Inches(0.58)))
        circ.fill.solid()
        circ.fill.fore_color.rgb = M2.GOLD
        circ.line.fill.background()
        circ.shadow.inherit = False
        tf = circ.text_frame
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = str(i + 1)
        run.font.size = Pt(25)
        run.font.bold = True
        run.font.color.rgb = ink
        run.font.name = "Calibri"
        # item titles follow the same title case as slide titles; the
        # one-line description underneath is a sentence, so it is left
        # alone
        rows = [([(M2._title_case(item[0].upper() + item[1:]),
                   {'bold': True, 'size': 25, 'color': ink})], 0,
                 {'bullet_style': 'none', 'space_before_pts': 0})]
        if i in hi:
            rows.append(([(desc, {'size': 22, 'color': M2.GRAY})], 0,
                         {'bullet_style': 'none', 'space_before_pts': 0}))
        # a shaded item renders its single line at the TOP of the
        # reserved two-row box, which sits high against the gold circle;
        # nudge it down to centre it (the current topic, which fills its
        # box with title + description, does NOT move)
        M2._add_hierarchical_bullets(
            slide, Inches(2.05), y if lit else int(y + M2.DIM_DROP),
            Inches(11.0), title_h + desc_h, rows, size=25,
            line_spacing_pts=0)
        # where this topic is taught — the pill dims with its row, so on a
        # section agenda only the current topic keeps its colour
        in_class = i in IN_CLASS_ITEMS
        if lit:
            pill_fill = M2.NAVY if in_class else M2.GOLD
            pill_ink = M2.WHITE if in_class else M2.NAVY
        else:
            pill_fill, pill_ink = M2.DIM, M2.WHITE
        pill_w = Inches(1.55)
        pill_y = y + Inches(0.02) + PILL_DROP
        M2._add_rounded_filled_box(
            slide, int(Inches(12.85) - pill_w), int(pill_y),
            int(pill_w), Inches(0.36), COVERAGE_LABEL[i],
            fill=pill_fill, text_color=pill_ink,
            size=13, corner_pct=0.30, shadow=lit)
        y = int(y + pitch)

    M2._draw_footer(slide, FOOTER, page_num)
    return slide


# --- package helpers -------------------------------------------------------

def slide_order(data):
    """display order -> slide part basenames."""
    pres = ET.fromstring(data["ppt/presentation.xml"])
    rid2t = {r.get("Id"): r.get("Target") for r in
             ET.fromstring(data["ppt/_rels/presentation.xml.rels"])}
    return [os.path.basename(rid2t[s.get(q(R, "id"))])
            for s in pres.find(q(P, "sldIdLst"))]


def slide_text(blob):
    return " ".join(e.text or "" for e in ET.fromstring(blob).iter(q(A, "t")))


def part_nums(data, folder, stem):
    pat = re.compile(r"^ppt/%s/%s(\d+)\.xml$" % (folder, stem))
    return [int(m.group(1)) for m in
            (pat.match(n) for n in data) if m]


def notes_part_for(data, slide_basename):
    """The notesSlide part a slide points at (or None)."""
    rels = "ppt/slides/_rels/%s.rels" % slide_basename
    if rels not in data:
        return None
    for r in ET.fromstring(data[rels]):
        if r.get("Type").endswith("/notesSlide"):
            return "ppt/notesSlides/" + os.path.basename(r.get("Target"))
    return None


def set_notes_text(blob, text):
    """Replace the body text of a notesSlide part with a single run."""
    tree = ET.fromstring(blob)
    for sp in tree.iter(q(P, "sp")):
        ph = sp.find(".//" + q(P, "nvSpPr") + "/" + q(P, "nvPr") + "/"
                     + q(P, "ph"))
        if ph is None or ph.get("type") != "body":
            continue
        tx = sp.find(q(P, "txBody"))
        for para in tx.findall(q(A, "p")):
            tx.remove(para)
        para = ET.SubElement(tx, q(A, "p"))
        run = ET.SubElement(para, q(A, "r"))
        t = ET.SubElement(run, q(A, "t"))
        t.text = text
        break
    return ET.tostring(tree, xml_declaration=True, encoding="UTF-8",
                       standalone=True)


def strip_creation_ids(blob):
    """Drop Office creationId GUIDs so a cloned part carries no
    duplicate ids."""
    tree = ET.fromstring(blob)
    for ext in list(tree.iter("{%s}ext" % A)) + list(tree.iter(q(P, "ext"))):
        if len(ext) and ext[0].tag.split("}")[-1] == "creationId":
            ext.getparent().remove(ext)
    for lst in list(tree.iter("{%s}extLst" % A)) + list(tree.iter(q(P, "extLst"))):
        if len(lst) == 0:
            lst.getparent().remove(lst)
    return ET.tostring(tree, xml_declaration=True, encoding="UTF-8",
                       standalone=True)


def insert_slide(data, after_display, sld_xml, notes_template, notes_text):
    """Add a brand-new slide (plus its notesSlide) right after
    *after_display* (1-based).  Returns the new slide's basename."""
    n = max(part_nums(data, "slides", "slide")) + 1
    nn = max(part_nums(data, "notesSlides", "notesSlide")) + 1
    spart = "ppt/slides/slide%d.xml" % n
    npart = "ppt/notesSlides/notesSlide%d.xml" % nn

    body = sld_xml.decode("utf-8")
    assert "r:id=" not in body and "r:embed=" not in body, \
        "generated slide references a relationship — rels needed"
    data[spart] = sld_xml
    data[npart] = set_notes_text(strip_creation_ids(notes_template),
                                 notes_text)

    def rels(pairs):
        root = ET.Element("{%s}Relationships" % PKG_R,
                          nsmap={None: PKG_R})
        for i, (typ, tgt) in enumerate(pairs, 1):
            ET.SubElement(root, "{%s}Relationship" % PKG_R, Id="rId%d" % i,
                          Type="http://schemas.openxmlformats.org/"
                               "officeDocument/2006/relationships/" + typ,
                          Target=tgt)
        return ET.tostring(root, xml_declaration=True, encoding="UTF-8",
                           standalone=True)

    data["ppt/slides/_rels/slide%d.xml.rels" % n] = rels([
        ("slideLayout", "../slideLayouts/slideLayout1.xml"),
        ("notesSlide", "../notesSlides/notesSlide%d.xml" % nn)])
    data["ppt/notesSlides/_rels/notesSlide%d.xml.rels" % nn] = rels([
        ("notesMaster", "../notesMasters/notesMaster1.xml"),
        ("slide", "../slides/slide%d.xml" % n)])

    # content types
    ct = ET.fromstring(data["[Content_Types].xml"])
    for pn, typ in ((spart, "slide"), (npart, "notesSlide")):
        ET.SubElement(ct, "{%s}Override" % CT, PartName="/" + pn,
                      ContentType="application/vnd.openxmlformats-"
                                  "officedocument.presentationml.%s+xml" % typ)
    data["[Content_Types].xml"] = ET.tostring(
        ct, xml_declaration=True, encoding="UTF-8", standalone=True)

    # presentation rels + sldIdLst
    prels = ET.fromstring(data["ppt/_rels/presentation.xml.rels"])
    rid = "rId%d" % (max(int(r.get("Id")[3:]) for r in prels) + 1)
    ET.SubElement(prels, "{%s}Relationship" % PKG_R, Id=rid,
                  Type="http://schemas.openxmlformats.org/officeDocument/"
                       "2006/relationships/slide",
                  Target="slides/slide%d.xml" % n)
    data["ppt/_rels/presentation.xml.rels"] = ET.tostring(
        prels, xml_declaration=True, encoding="UTF-8", standalone=True)

    pres = ET.fromstring(data["ppt/presentation.xml"])
    lst = pres.find(q(P, "sldIdLst"))
    new_id = max(int(s.get("id")) for s in lst) + 1
    el = ET.Element(q(P, "sldId"))
    el.set("id", str(new_id))
    el.set(q(R, "id"), rid)
    lst.insert(after_display, el)
    data["ppt/presentation.xml"] = ET.tostring(
        pres, xml_declaration=True, encoding="UTF-8", standalone=True)
    return "slide%d.xml" % n


def topbar_tag_element(tree):
    """The <a:t> holding a slide's top-bar section tag, or None. Detected
    by GEOMETRY (a top-level text shape sitting in the 0.35" top-bar band)
    plus the "Module 3 · " prefix — never by shape name."""
    spTree = tree.find(q(P, "cSld") + "/" + q(P, "spTree"))
    for sp in spTree.findall(q(P, "sp")):
        off = sp.find(q(P, "spPr") + "/" + q(A, "xfrm") + "/" + q(A, "off"))
        if off is None or int(off.get("y")) > 320000:
            continue
        tx = sp.find(q(P, "txBody"))
        if tx is None:
            continue
        for para in tx.findall(q(A, "p")):
            runs = para.findall(q(A, "r"))
            if not runs:
                continue
            t = runs[0].find(q(A, "t"))
            if t is None or not (t.text or "").startswith(TOPBAR_PREFIX):
                continue
            for extra in runs[1:]:          # tag is one run by convention
                para.remove(extra)
            return t
    return None


def retag_top_bars(data, spec_by_display):
    """Set every content slide's section tag to the agenda item it sits
    under. Slides ahead of the first agenda slide keep their own tag, as
    do TOPBAR_KEEP tags and slides with no top bar (title, polls,
    backups)."""
    cur = None
    changed = []
    for i, base in enumerate(slide_order(data), 1):
        part = "ppt/slides/" + base
        tree = ET.fromstring(data[part])
        t = topbar_tag_element(tree)
        if i in spec_by_display:
            hl = spec_by_display[i].get("highlight")
            cur = M3_OUTLINE[min(hl)][0] if hl else None
            new = TAG
        elif t is None or t.text in TOPBAR_KEEP:
            continue
        elif cur is None:                    # front matter
            new = FRONT_TAG_FIX.get(t.text, t.text)
        else:
            new = TOPBAR_PREFIX + cur
        if t is None or t.text == new:
            continue
        old, t.text = t.text, new
        data[part] = ET.tostring(tree, xml_declaration=True,
                                 encoding="UTF-8", standalone=True)
        changed.append((i, old, new))
    return changed


def refresh_cached_page_numbers(data):
    """Footer page numbers are live <a:fld type="slidenum"> fields; keep
    the CACHED text in step with the new display order so the deck looks
    right before PowerPoint recomputes."""
    for i, base in enumerate(slide_order(data), 1):
        part = "ppt/slides/" + base
        tree = ET.fromstring(data[part])
        touched = False
        for fld in tree.iter(q(A, "fld")):
            if fld.get("type") != "slidenum":
                continue
            t = fld.find(q(A, "t"))
            if t is not None and t.text != str(i):
                t.text = str(i)
                touched = True
        if touched:
            data[part] = ET.tostring(tree, xml_declaration=True,
                                     encoding="UTF-8", standalone=True)


# --- read the deck, locate the agenda slides ------------------------------
z = zipfile.ZipFile(DECK)
data = {n: z.read(n) for n in z.namelist()}
z.close()
order = slide_order(data)
agenda = [i + 1 for i, base in enumerate(order)
          if TITLE in slide_text(data["ppt/slides/" + base])]
print("agenda slides found at displays %s" % agenda)

inserted_at = None
if len(agenda) == len(SPECS) - 1:
    anchors = [i + 1 for i, base in enumerate(order)
               if INSERT_AFTER in slide_text(data["ppt/slides/" + base])]
    if len(anchors) != 1:
        raise SystemExit("anchor %r matches %d slides, need exactly 1"
                         % (INSERT_AFTER, len(anchors)))
    inserted_at = anchors[0] + 1
    print("agenda slide for %r is missing — will insert it at display %d"
          % (M3_OUTLINE[min(SPECS[INSERT_SPEC]['highlight'])][0],
             inserted_at))
elif len(agenda) != len(SPECS):
    raise SystemExit("expected %d or %d agenda slides, found %d"
                     % (len(SPECS) - 1, len(SPECS), len(agenda)))

# final display numbers of the agenda slides
finals = list(agenda)
if inserted_at is not None:
    finals = [d + 1 if d >= inserted_at else d for d in finals]
    finals.insert(INSERT_SPEC, inserted_at)

# --- generate the replacement slides --------------------------------------
prs = M2.Presentation()
prs.slide_width = int(M2.SLIDE_W)
prs.slide_height = int(M2.SLIDE_H)
for spec, disp in zip(SPECS, finals):
    make_m3_outline(prs, disp, descriptions=spec.get("descriptions", False),
                    highlight_set=spec.get("highlight"))
tmp = HERE / "_agenda_tmp.pptx"
prs.save(str(tmp))
z = zipfile.ZipFile(tmp)
tdata = {n: z.read(n) for n in z.namelist()}
z.close()
torder = slide_order(tdata)
print("generated %d replacement slides" % len(SPECS))


def gen_sptree(idx):
    ttree = ET.fromstring(tdata["ppt/slides/" + torder[idx]])
    return ttree.find(".//" + q(P, "cSld") + "/" + q(P, "spTree"))


# --- insert the missing agenda slide (first run only) ----------------------
if inserted_at is not None:
    tmpl = notes_part_for(data, order[agenda[0] - 1])   # any agenda notes
    base = insert_slide(data, inserted_at - 1,
                        tdata["ppt/slides/" + torder[INSERT_SPEC]],
                        data[tmpl], INSERT_NOTES)
    print("inserted %s at display %d" % (base, inserted_at))
    order = slide_order(data)

# --- transplant spTrees ----------------------------------------------------
for idx, disp in enumerate(finals):
    part = "ppt/slides/" + order[disp - 1]
    tree = ET.fromstring(data[part])
    csld = tree.find(q(P, "cSld"))
    csld.replace(csld.find(q(P, "spTree")), gen_sptree(idx))
    timing = tree.find(q(P, "timing"))
    if timing is not None:
        tree.remove(timing)
    data[part] = ET.tostring(tree, xml_declaration=True,
                             encoding="UTF-8", standalone=True)
    spec = SPECS[idx]
    print("display %2d (%s) <- %s"
          % (disp, part, "overview (all items)" if spec.get("descriptions")
             else "items %s" % sorted(x + 1 for x in spec["highlight"])))

for disp, old, new in retag_top_bars(data, dict(zip(finals, SPECS))):
    print("tag %2d: %-38s -> %s" % (disp, old, new))

refresh_cached_page_numbers(data)

out = DECK.with_suffix(".retrofit_tmp.pptx")
with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as zout:
    for name, blob in data.items():
        zout.writestr(name, blob)
shutil.move(str(out), str(DECK))
tmp.unlink()
print("saved %s (%d slides)" % (DECK.name, len(slide_order(data))))
