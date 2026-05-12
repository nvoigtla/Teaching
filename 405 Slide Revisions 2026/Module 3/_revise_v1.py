"""
Module 3 revision v1: bring the deck to 2026 currency.

  1. Tesla → Rivian on slides 31, 35, 36, 37, 39, 40, 54, 55, 56
     (long-run input mix + cost function).
  2. Designer → AI researcher on slides 26, 28 (and update the math: $5M / $3.5M / $8M MC).
  3. Burn60 → ChatGPT subscription tiers on slides 50, 52.
  4. iPhone 11 Pro Max → iPhone 17 on slide 59 (with $1,199 / $580 AVC).
  5. Airbus A380/A318 → Alphabet on slide 70 (economies of scope).
  6. Insert NEW slide after position 46: Meta Reality Labs (modern sunk-cost parallel).
  7. Insert NEW slide between positions 67 and 68: AI training costs (economies of scale).
  8. After both inserts, walk all slides and update the chrome page numbers.
  9. Replace every slide's speaker note from the updated NOTES dict.

Direct zip + lxml surgery throughout. python-pptx is used only on
clean sidecar decks (no risk of NULL-rel stripping).
"""

import sys
import zipfile
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches

HERE = Path(__file__).parent
sys.path.insert(0, str(HERE))
from _build_template_samples import (  # noqa: E402
    NAVY, GRAY, RULE, GOLD, WHITE,
    MARGIN, RULE_W, GOLD_W, SLIDE_W, SLIDE_H, FOOTER_TEXT,
    _add_text, _add_rect,
    _draw_top_bar, _draw_action_title, _draw_footer,
    _add_bulleted_list,
)
from _speaker_notes import NOTES  # noqa: E402
from _polish_v2 import (  # noqa: E402
    NS, Q, parse, serialize, _resolve_relative,
    REL_SLIDE, REL_LAYOUT, REL_IMAGE, REL_NOTES, REL_NOTES_MASTER,
    CT_NOTES,
    get_slide_path_at_position, next_rId,
    build_notes_slide_xml,
    add_notes_to_slide,
)

TARGET = HERE / "Module 3_NEW_draft.pptx"
SIDECAR_META = HERE / "_sidecar_meta.pptx"
SIDECAR_AI_SCALE = HERE / "_sidecar_ai_scale.pptx"


# ============================================================
# Helpers
# ============================================================

def text_replace_in_slide(members, slide_path, replacements):
    """Replace text in all <a:t> elements within the slide.

    *replacements* is an ordered list of (old, new) tuples. Order
    matters: longer phrases should come first.
    """
    slide = parse(members[slide_path])
    n = 0
    for t in slide.iter(Q('a', 't')):
        if not t.text:
            continue
        original = t.text
        new_text = original
        for old, new in replacements:
            if old in new_text:
                new_text = new_text.replace(old, new)
        if new_text != original:
            t.text = new_text
            n += 1
    if n > 0:
        members[slide_path] = serialize(slide)
    return n


def _strip_slide(slide):
    spTree = slide.shapes._spTree
    for child in list(spTree):
        tag = child.tag.split('}', 1)[1]
        if tag in ('sp', 'grpSp', 'pic', 'graphicFrame', 'cxnSp', 'contentPart'):
            spTree.remove(child)
    slide._element.set('showMasterSp', '0')


def _new_blank_slide(prs):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _strip_slide(slide)
    return slide


def next_slide_filename(members):
    used = set()
    for name in members:
        if name.startswith('ppt/slides/slide') and name.endswith('.xml'):
            tail = name[len('ppt/slides/slide'):-len('.xml')]
            try:
                used.add(int(tail))
            except ValueError:
                pass
    n = 1
    while n in used:
        n += 1
    return f'ppt/slides/slide{n}.xml'


# ============================================================
# (1) Tesla → Rivian on long-run + cost-function slides
# ============================================================

RIVIAN_REPLACEMENTS = [
    # Longer phrases first
    ("Tesla’s New Gigafactory", "Rivian's New Georgia plant"),
    ("Tesla's New Gigafactory", "Rivian's New Georgia plant"),
    ("Tesla’s Gigafactory", "Rivian's Georgia plant"),
    ("Tesla's Gigafactory", "Rivian's Georgia plant"),
    ("Tesla's Current Plan", "Rivian's Current Plan"),
    ("Tesla’s Current Plan", "Rivian's Current Plan"),
    ("Tesla builds a new Gigafactory", "Rivian builds its new Georgia plant"),
    ("Tesla is building a new Gigafactory in Berlin, Germany",
     "Rivian is building a new plant in Stanton Springs, Georgia"),
    ("Tesla is planning to produce 500 cars per week",
     "Rivian is planning to produce 500 vehicles per week"),
    ("500 cars per week", "500 vehicles per week"),
    ("Tesla's input mix", "Rivian's input mix"),
    ("Tesla’s input mix", "Rivian's input mix"),
    ("Berlin Gigafactory", "Georgia plant"),
    ("new Gigafactory", "new Georgia plant"),
    ("Gigafactory", "Georgia plant"),
    ("Tesla", "Rivian"),
    # Production-volume terminology
    ("of cars", "of vehicles"),
    ("of Cars", "of Vehicles"),
]

RIVIAN_SLIDES = (31, 35, 36, 37, 39, 40, 54, 55, 56)


def apply_rivian_swap(members):
    print("[1] Tesla → Rivian:")
    for pos in RIVIAN_SLIDES:
        slide_path = get_slide_path_at_position(members, pos)
        n = text_replace_in_slide(members, slide_path, RIVIAN_REPLACEMENTS)
        print(f"    slide {pos} ({slide_path}): {n} <a:t> elements updated")


# ============================================================
# (2) Designer → AI researcher on slide 26, 28
# ============================================================

DESIGNER_TO_AI_REPLACEMENTS = [
    # Slide 26 specifics
    ("The full cost of a new Tesla Designer",
     "The full cost of poaching an AI researcher"),
    ("Tesla is trying to hire a star designer from BMW for its next models",
     "Anthropic is trying to poach a star researcher from Google DeepMind"),
    ("Designer says she’d be willing to join for $2M annual salary",
     "Researcher says she'd be willing to join for $5M annual salary"),
    ("Designer says she'd be willing to join for $2M annual salary",
     "Researcher says she'd be willing to join for $5M annual salary"),
    ("Tesla already employs 2 star designers, each earning $1.5M",
     "Anthropic already employs 2 star researchers, each earning $3.5M"),
    ("If the new designer is hired, the two existing ones would demand the same salary",
     "If the new researcher is hired, the two existing ones would demand the same salary"),
    ("What is marginal cost of the 3rd designer?",
     "What is the marginal cost of the 3rd researcher?"),
    # Slide 28 specifics
    ("The star designer herself will have to be paid $2M",
     "The star researcher herself will have to be paid $5M"),
    ("The two existing designers would earn $0.5M extra",
     "The two existing researchers would earn $1.5M extra each"),
    ("The marginal cost of the 3rd designer is thus $2M+2*$0.5M = $3M",
     "The marginal cost of the 3rd researcher is thus $5M+2*$1.5M = $8M"),
    # Revised-title overlay (added by chrome step earlier)
    ("Solution: marginal cost of the 3rd designer = $3M",
     "Solution: marginal cost of the 3rd researcher = $8M"),
]


def apply_ai_researcher_swap(members):
    print("[2] Designer → AI researcher:")
    for pos in (26, 28):
        slide_path = get_slide_path_at_position(members, pos)
        n = text_replace_in_slide(members, slide_path, DESIGNER_TO_AI_REPLACEMENTS)
        print(f"    slide {pos}: {n} <a:t> elements updated")


# ============================================================
# (3) Burn60 → ChatGPT subscription tiers on slides 50, 52
# ============================================================

CHATGPT_REPLACEMENTS = [
    # Slide 50 specifics
    ("Marginal Cost in Action:", "Marginal Cost in Action:"),
    ("What is the MC of a Burn60 Workout?",
     "What's the MC of adding a 2nd user to ChatGPT?"),
    ("PACKAGES", "OPENAI TIERS"),
    ("10 CLASSES FOR $10 EACH", "PLUS: $20 / user / month"),
    ("20 CLASSES for $8 EACH", "TEAM: $25 / user / month, 2 users min."),
    ("Buy Now", "Subscribe"),
    # Slide 52 specifics
    ("TC of 20 lessons: 20*$8=$160",
     "TC of 2 users on Team: 2 × $25 = $50/month"),
    ("TC of 10 lessons: 10*$10=$100",
     "TC of 1 user on Plus: 1 × $20 = $20/month"),
    ("So 10 additional sessions cost $160-$100=$60, or $6 per session.",
     "So the 2nd user costs $50 − $20 = $30/month – more than Team's $25 sticker."),
    ("More significant discount than it seems at first.",
     "MC > AC: subscription tiers can hide a higher marginal cost."),
    ("Intuition: “save” $2 on the first 10 sessions as well when buying 20 at $8.",
     "Intuition: you re-price the existing user ($20 → $25) and add a new one ($25)."),
    ("Intuition: \"save\" $2 on the first 10 sessions as well when buying 20 at $8.",
     "Intuition: you re-price the existing user ($20 → $25) and add a new one ($25)."),
    # Revised-title overlay (slide 50)
    ("Marginal cost ≠ average cost: Burn60 workout packages",
     "Marginal cost ≠ average cost: ChatGPT subscription tiers"),
    # Revised-title overlay (slide 52)
    ("Solution: MC = $6", "Solution: MC = $30/user · month"),
]


def apply_chatgpt_swap(members):
    print("[3] Burn60 → ChatGPT:")
    for pos in (50, 52):
        slide_path = get_slide_path_at_position(members, pos)
        n = text_replace_in_slide(members, slide_path, CHATGPT_REPLACEMENTS)
        print(f"    slide {pos}: {n} <a:t> elements updated")


# ============================================================
# (4) iPhone 11 Pro Max → iPhone 17 on slide 59
# ============================================================

IPHONE_REPLACEMENTS = [
    ("AVC of iphone 11 Pro Max", "AVC of iphone 17"),
    ("AVC of iPhone 11 Pro Max", "AVC of iPhone 17"),
    ("iphone 11 Pro Max", "iphone 17"),
    ("iPhone 11 Pro Max", "iPhone 17"),
    ("Retail price: $1099", "Retail price: $1,199"),
    ("TVC: ~$500", "TVC: ~$580"),
]


def apply_iphone_update(members):
    print("[4] iPhone 11 → iPhone 17:")
    slide_path = get_slide_path_at_position(members, 59)
    n = text_replace_in_slide(members, slide_path, IPHONE_REPLACEMENTS)
    print(f"    slide 59: {n} <a:t> elements updated")


# ============================================================
# (5) Airbus A380/A318 → Alphabet on slide 70
# ============================================================

ALPHABET_REPLACEMENTS = [
    ("Example: Airbus A380 and A318",
     "Example: Alphabet (Google)"),
    ("Airbus A380 and A318",
     "Alphabet (Google)"),
    # Revised-title overlay
    ("Sharing capabilities across products: economies of scope",
     "Sharing capabilities across products: economies of scope"),
]

# Slide 70 bullets ("Input production, Engineering know-how, R&D,
# Marketing") apply to Alphabet too – no replacement needed there.


def apply_alphabet_swap(members):
    print("[5] Airbus → Alphabet:")
    slide_path = get_slide_path_at_position(members, 70)
    n = text_replace_in_slide(members, slide_path, ALPHABET_REPLACEMENTS)
    print(f"    slide 70: {n} <a:t> elements updated")


# ============================================================
# (6) NEW slide – Meta Reality Labs (after position 46)
# ============================================================

def build_meta_sidecar():
    prs = Presentation()
    prs.slide_width = Emu(int(SLIDE_W))
    prs.slide_height = Emu(int(SLIDE_H))
    slide = _new_blank_slide(prs)

    _draw_top_bar(slide, "Module 3 · Part 2 · Costs")
    _draw_action_title(slide,
        "Modern sunk cost: Meta's Reality Labs has lost $50B+ since 2020")

    bullets = [
        "Meta has poured ~$50B into Reality Labs (Metaverse, VR, AR) since 2020",
        "Wall Street keeps asking when it pays off – Zuckerberg keeps investing",
        "Past losses are sunk – ignore them when deciding what to do next",
        "Right question: does the next $10B have positive expected value going forward?",
        "Same lesson as Waterworld – classic sunk-cost discipline, 2020s edition",
    ]
    _add_bulleted_list(
        slide,
        left=MARGIN, top=Inches(1.85),
        width=RULE_W, height=Inches(5.0),
        items=bullets,
        size=22, color=NAVY, bullet_color=NAVY,
        bullet_char="▪",
        bullet_size_pct=150,
        line_spacing_pts=20,
    )

    _draw_footer(slide, FOOTER_TEXT, 47)   # page number set later in pass 8
    prs.save(SIDECAR_META)


# ============================================================
# (7) NEW slide – AI training costs (between current 67 and 68)
# ============================================================

def build_ai_scale_sidecar():
    prs = Presentation()
    prs.slide_width = Emu(int(SLIDE_W))
    prs.slide_height = Emu(int(SLIDE_H))
    slide = _new_blank_slide(prs)

    _draw_top_bar(slide, "Module 3 · Part 2 · Costs")
    _draw_action_title(slide,
        "AI training: classic economies of scale at extreme cost")

    bullets = [
        "Training a frontier model (GPT-5, Claude Opus 4.7) costs $500M+ in compute",
        "Once trained, marginal cost per query is fractions of a cent",
        "Spread fixed cost across 300M+ users → near-zero per-user cost",
        "NVIDIA H100/B200 GPU clusters cost $40M+ each – lumpy capital",
        "Why only a handful of labs compete in foundation models",
    ]
    _add_bulleted_list(
        slide,
        left=MARGIN, top=Inches(1.85),
        width=RULE_W, height=Inches(5.0),
        items=bullets,
        size=22, color=NAVY, bullet_color=NAVY,
        bullet_char="▪",
        bullet_size_pct=150,
        line_spacing_pts=20,
    )

    _draw_footer(slide, FOOTER_TEXT, 69)   # page number set later
    prs.save(SIDECAR_AI_SCALE)


def inject_new_slide(members, sidecar_path, insert_after_pos):
    """Build new slide from sidecar and inject it into the deck just
    AFTER position *insert_after_pos* (1-based, in the current deck)."""
    with zipfile.ZipFile(sidecar_path, 'r') as zf:
        new_slide_xml = zf.read('ppt/slides/slide1.xml')

    # Pick a unique slide filename
    new_slide_path = next_slide_filename(members)
    slide_dir, slide_file = new_slide_path.rsplit('/', 1)
    new_rels_path = f'{slide_dir}/_rels/{slide_file}.rels'

    members[new_slide_path] = new_slide_xml

    # Minimal rels: just layout
    rels_root = etree.Element(Q('rel', 'Relationships'),
                              nsmap={None: NS['rel']})
    layout_rel = etree.SubElement(rels_root, Q('rel', 'Relationship'))
    layout_rel.set('Id', 'rId1')
    layout_rel.set('Type', REL_LAYOUT)
    layout_rel.set('Target', '../slideLayouts/slideLayout2.xml')
    members[new_rels_path] = serialize(rels_root)

    # Content_Types Override
    ct = parse(members['[Content_Types].xml'])
    ov = etree.SubElement(ct, Q('ct', 'Override'))
    ov.set('PartName', '/' + new_slide_path)
    ov.set('ContentType',
           'application/vnd.openxmlformats-officedocument.presentationml.slide+xml')
    members['[Content_Types].xml'] = serialize(ct)

    # Add presentation-level Relationship + sldIdLst entry
    pres = parse(members['ppt/presentation.xml'])
    pres_rels = parse(members['ppt/_rels/presentation.xml.rels'])
    new_rId = next_rId(pres_rels)
    new_rel = etree.SubElement(pres_rels, Q('rel', 'Relationship'))
    new_rel.set('Id', new_rId)
    new_rel.set('Type', REL_SLIDE)
    new_rel.set('Target', new_slide_path[len('ppt/'):])
    members['ppt/_rels/presentation.xml.rels'] = serialize(pres_rels)

    sldIdLst = pres.find(Q('p', 'sldIdLst'))
    existing_ids = [int(s.get('id'))
                    for s in sldIdLst.findall(Q('p', 'sldId'))]
    new_id = max(existing_ids) + 1 if existing_ids else 256
    new_sld = etree.Element(Q('p', 'sldId'))
    new_sld.set('id', str(new_id))
    new_sld.set(Q('r', 'id'), new_rId)
    # Insert after the slide at *insert_after_pos* (so it becomes pos+1)
    sldIdLst.insert(insert_after_pos, new_sld)   # 0-based index = insert_after_pos
    members['ppt/presentation.xml'] = serialize(pres)

    print(f"    injected new slide {new_slide_path} at position {insert_after_pos + 1}")


# ============================================================
# (8) Update chrome page numbers after structural changes
# ============================================================

def update_page_numbers(members):
    print("[8] Update chrome page numbers:")
    pres = parse(members['ppt/presentation.xml'])
    pres_rels = parse(members['ppt/_rels/presentation.xml.rels'])
    rId_to_target = {r.get('Id'): ('ppt/' + r.get('Target')).replace('\\', '/')
                     for r in pres_rels.findall(Q('rel', 'Relationship'))
                     if r.get('Type') == REL_SLIDE}
    sld_seq = pres.find(Q('p', 'sldIdLst')).findall(Q('p', 'sldId'))

    n_updated = 0
    for i, sld in enumerate(sld_seq, start=1):
        rId = sld.get(Q('r', 'id'))
        slide_path = rId_to_target[rId]
        # The chrome page number is the textbox at (12.6", 7.22"),
        # size (0.5", 0.3"), containing a digit-only string.
        slide = parse(members[slide_path])
        for sp in slide.iter(Q('p', 'sp')):
            xfrm = sp.find(Q('p', 'spPr') + '/' + Q('a', 'xfrm'))
            if xfrm is None:
                continue
            off = xfrm.find(Q('a', 'off'))
            ext = xfrm.find(Q('a', 'ext'))
            if off is None or ext is None:
                continue
            # Look for textbox near (12.6", 7.22") with size (~0.5", ~0.3")
            ox, oy = int(off.get('x')), int(off.get('y'))
            cx, cy = int(ext.get('cx')), int(ext.get('cy'))
            EMU_PER_IN = 914400
            if (abs(ox - int(12.6 * EMU_PER_IN)) > 50000 or
                abs(oy - int(7.22 * EMU_PER_IN)) > 50000):
                continue
            if (abs(cx - int(0.5 * EMU_PER_IN)) > 30000 or
                abs(cy - int(0.3 * EMU_PER_IN)) > 30000):
                continue
            # Update the text inside this sp
            for t in sp.iter(Q('a', 't')):
                if t.text and t.text.strip().isdigit():
                    if t.text != str(i):
                        t.text = str(i)
                        n_updated += 1
                    break
        members[slide_path] = serialize(slide)
    print(f"    {n_updated} page numbers updated")


# ============================================================
# (9) Refresh ALL speaker notes from the updated NOTES dict
# ============================================================

def refresh_all_notes(members):
    print("[9] Refresh speaker notes for all slides:")
    pres = parse(members['ppt/presentation.xml'])
    n_total = len(pres.find(Q('p', 'sldIdLst')).findall(Q('p', 'sldId')))
    if n_total != len(NOTES):
        print(f"    WARN: deck has {n_total} slides but NOTES has {len(NOTES)}")
    for pos in range(1, n_total + 1):
        if pos in NOTES:
            add_notes_to_slide(members, pos, NOTES[pos])
    print(f"    notes refreshed for {n_total} slides")


# ============================================================
# Main
# ============================================================

def main():
    print(f"Reading {TARGET.name}")
    with zipfile.ZipFile(TARGET, 'r') as zf:
        members = {n: zf.read(n) for n in zf.namelist()}

    apply_rivian_swap(members)
    apply_ai_researcher_swap(members)
    apply_chatgpt_swap(members)
    apply_iphone_update(members)
    apply_alphabet_swap(members)

    # ---- Inserts (must happen AFTER position-indexed text edits) ----
    print("[6] Insert NEW Meta Reality Labs slide (after position 46):")
    build_meta_sidecar()
    inject_new_slide(members, SIDECAR_META, insert_after_pos=46)

    print("[7] Insert NEW AI training-cost slide (between 67 and 68):")
    build_ai_scale_sidecar()
    # After Meta insert, current 67 became new 68. We want the AI slide
    # AFTER what is now position 68 (i.e., the old slide 67 "Technological
    # Reasons for Economies of Scale"), so insert after position 68.
    inject_new_slide(members, SIDECAR_AI_SCALE, insert_after_pos=68)

    update_page_numbers(members)
    refresh_all_notes(members)

    # Write atomically
    tmp = TARGET.with_suffix('.pptx.tmp')
    with zipfile.ZipFile(tmp, 'w', zipfile.ZIP_DEFLATED) as zf:
        for n, d in members.items():
            zf.writestr(n, d)
    tmp.replace(TARGET)
    print(f"\nWrote {TARGET.name}")


if __name__ == "__main__":
    main()
