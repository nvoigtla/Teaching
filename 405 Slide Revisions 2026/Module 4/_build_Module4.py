# ###########################################################################
#  STALE AS OF 2026-08-30 - DO NOT RUN.
#
#  Module 4 is finished and Nico is hand-editing "Module 4 - Revised.pptx"
#  directly.  The .pptx is now the source of truth; running this script
#  would regenerate the deck from scratch and destroy every hand edit made
#  since.  Kept for reference - it records how each slide was built and
#  carries the dated comments for every hand-edit port up to the freeze.
#
#  To change a slide now: edit the .pptx in place with zip + lxml surgery.
#  NEVER round-trip it through python-pptx - that strips the PollEverywhere
#  "tags" parts and NULL rels, and this deck has 7 live poll slides, any one
#  of which will crash the full-screen slideshow DECK-WIDE if broken.
#
#  _animate.py is still safe to run (it is OOXML surgery, not a round-trip),
#  but its plans are keyed by DISPLAY NUMBER - renumber them if slides are
#  inserted or deleted by hand.
# ###########################################################################

# ==========================================================================
#  _build_Module4.py — build script for "Module 4 - Revised.pptx"
#
#  Competitive Markets and Market Interventions, 83 slides per the approved
#  outline in "Module 4 - Revised - outline.md".  The six PollEverywhere
#  slides are positional stubs until the _splice_media.py pass.
#
#  Primitives come from _m4_helpers.py (carved out of Module 1's build
#  script by _make_helpers.py); speaker notes of the original deck come
#  verbatim from _m4_notes.py.  This script is the SOURCE OF TRUTH for the
#  deck — hand edits in PowerPoint are ported back into it.
#
#  Pipeline (rerunnable, Module 7 pattern):
#      python _build_Module4.py            -> Module 4 - Revised.pptx
#      python _splice_media.py             -> polls, verbatim + notes + tags
#      python _group_pass.py               -> box+text / figure+shade groups
#      python _animate.py all apply        -> fade builds
# ==========================================================================

from pathlib import Path

from pptx import Presentation

import _m4_helpers as _H
from _m4_helpers import *                                        # noqa: F401,F403
from _m4_helpers import (
    CREAM, DARKRED, DIM, GOLD, GOLD_W, GRAY, MARGIN, NAVY, RED, RULE,
    RULE_W, SLIDE_H, SLIDE_W, WHITE, FADED, BLUE_PED,
    Inches, Pt, PP_ALIGN, MSO_ANCHOR, MSO_SHAPE, RGBColor, qn,
    _add_arrow, _add_arrow_shape, _add_convention_box, _add_drop_shadow,
    _add_graphicframe_shadow, _add_hierarchical_bullets, _add_math_equation,
    _add_media_image, _add_mixed_textbox, _add_outlined_box,
    _add_pollbreak_badge, _add_ps_pointer, _add_rect,
    _add_rounded_filled_box, _add_slidenum_field, _add_styled_table,
    _add_takeaway_bar, _add_text, _apply_picture_style, _blank_slide,
    _draw_action_title, _draw_footer, _draw_poll_pill, _draw_top_bar_tc,
    _fig_axes, _fig_curve_label, _fig_guide, _fig_line, _fig_xlab, _fig_ylab,
    _fig_underbrace, _fig_vbrace, _welfare_rows, GREEN_DK, _text_w_in,
    _fig_ylab_subsup, EMC_PURPLE,
    _omml_frac, _omml_run, _omml_sub, _omml_sup, _omml_text,
    _omml_underbrace,
    _set_fill_alpha, _set_notes, _title_case, apply_symbol_subscripts,
    content_slide, make_diagram_slide, make_stub, SimpleFig,
)
from _m4_notes import NOTES

# Two of my original notes quote the OLD figures, so they would contradict
# the slide after the 2026-08-28 switch.  They are overridden here rather
# than edited in _m4_notes.py, which stays a verbatim capture.
NOTES = dict(NOTES)
NOTES[41] = (
    "A price taker\u2019s supply curve. We can trace out the supply curve "
    "of the firm.\nAt P = 210, you maximize profits by producing 187.5 "
    "tons.\nAt P = 400, you expand more.\nBottom line: the MC curve is "
    "the firm\u2019s short-run supply curve. We can draw the supply curve."
)
NOTES[19] = (
    "Market price: $400 per ton is about $20 per cwt, which is the USDA "
    "season-average price for fresh-market cabbage in 2020. Recent years "
    "run higher \u2014 $31.10 per cwt in 2024.\n"
    "https://www.agmrc.org/commodities-products/vegetables/cabbage"
)

OUT = Path(__file__).parent / "Module 4 - Revised.pptx"

# 2026-08-29 (Nico): the course-layer "dark yellow" (#B8860B) — used on
# slide 10 for the market demand curve and its D label, where plain GOLD
# washes out against white.
DARKYELLOW = RGBColor(0xB8, 0x86, 0x0B)

# --------------------------------------------------------------------------
#  The OMML helpers build raw XML by string concatenation and do NOT escape
#  their argument, so a literal "<" in a comparison ("P < AVC") silently
#  produces an invalid slide part.  Module 4 is full of such comparisons, so
#  the two text-producing helpers are wrapped here rather than fixing them
#  at every call site.
# --------------------------------------------------------------------------

_omml_text_raw = _omml_text
_omml_run_raw = _omml_run


def _xml_escape(text):
    return (text.replace('&', '&amp;').replace('<', '&lt;')
                .replace('>', '&gt;'))


def _omml_text(text, **kw):                                  # noqa: F811
    return _omml_text_raw(_xml_escape(text), **kw)


def _omml_run(text, **kw):                                   # noqa: F811
    return _omml_run_raw(_xml_escape(text), **kw)

# --------------------------------------------------------------------------
#  Deck constants
# --------------------------------------------------------------------------

_H.FOOTER_TEXT = ("Management 405  ·  Module 4  ·  "
                  "Competitive Markets and Market Interventions")
FOOTER_TEXT = _H.FOOTER_TEXT

DECK_TITLE = "Competitive Markets and Market Interventions"
DECK_SUB = "Module 4"

# Top-bar tags.  One constant per outline item, derived from the item
# titles below, so renaming an item renames every tag that uses it.
# 2026-08-30: the module front matter sits INSIDE the introduction
# video's block, so it carries that video's number (Teaching
# CLAUDE.md, top-bar tag rule).  The outline slides are that video's
# agenda, so they read "Video 1 · Agenda".
TAG_LOG = "Module 4 · Video 1 · Logistics"
TAG_ROADMAP = "Module 4 · Video 1 · Course Roadmap"
TAG_OUTLINE = "Module 4 · Video 1 · Agenda"
TAG_SUMMARY = "Module 4 · Summary"

# The poll slides are PLACEHOLDERS for now (2026-08-28, Nico): the two
# cabbage activities still offer the retired answers, so the live slides
# are deliberately not spliced in.  Each placeholder title carries
# "++UPDATE" so it is obvious at a glance which slides are still pending.
POLL_MARK = "++UPDATE"
STUB_POLL = ("PollEverywhere placeholder — re-key the activity in PollEv, "
             "then run _splice_media.py")

# --------------------------------------------------------------------------
#  The Yi-family cost function.
#
#  These five constants drive the cost table, the TC chart, both worked
#  solutions, the profit rectangle and the supply-curve traces.  EVERY
#  printed figure below is derived from them through _num() and the
#  TC/MC/AVC omml builders — nothing is typed twice.
#
#  2026-08-28 (Nico): switched from my original 30,000 + 40Q + 0.2Q² at
#  P = 230/160 to these figures, which are the ones in MW's deck.  They are
#  the more defensible pair: ATC at Q* is $358/ton against the $348/ton in
#  the UC Ventura County cost study my speaker notes already cite, and a
#  $400/ton price is USDA's 2020 season average.  See
#  _cabbage_numbers_check.md.  The two PollEverywhere activities are keyed
#  to the OLD answers and have to be re-entered in the PollEv account.
# --------------------------------------------------------------------------

TFC = 60000          # total fixed cost, $
B_LIN = 135          # linear term of TC, $ per ton
B_QUAD = 0.2         # quadratic term of TC
P_HIGH = 400         # first market price, $ per ton
P_LOW = 210          # market price after the Chinese imports, $ per ton


def tc(q):
    return TFC + B_LIN * q + B_QUAD * q * q


def mc(q):
    return B_LIN + 2 * B_QUAD * q


def avc(q):
    return B_LIN + B_QUAD * q


def q_star(p):
    return (p - B_LIN) / (2 * B_QUAD)


Q_HIGH = q_star(P_HIGH)          # 662.5
Q_LOW = q_star(P_LOW)            # 187.5


def _num(x, dp=None):
    """Money / quantity as the slides print it: thousands separated, and
    only as many decimals as the value actually has (max 2)."""
    if dp is None:
        if abs(x - round(x)) < 1e-9:
            dp = 0
        elif abs(x * 10 - round(x * 10)) < 1e-9:
            dp = 1
        else:
            dp = 2
    return "{:,.{}f}".format(x, dp)


def _omml_tc():
    """TC = F + bQ + cQ²  as OMML."""
    return (_omml_text('TC') + _omml_text(' = %s + %s' % (_num(TFC),
                                                          _num(B_LIN)))
            + _omml_run('Q') + _omml_text(' + %s' % _num(B_QUAD))
            + _omml_sup(_omml_run('Q'), _omml_text('2')))


def _omml_mc_expr():
    """b + 2cQ  as OMML (the marginal-cost expression, no left-hand side)."""
    return (_omml_text('%s + %s' % (_num(B_LIN), _num(2 * B_QUAD)))
            + _omml_run('Q'))


def _omml_avc_expr():
    """b + cQ  as OMML."""
    return (_omml_text('%s + %s' % (_num(B_LIN), _num(B_QUAD)))
            + _omml_run('Q'))

# --------------------------------------------------------------------------
#  Module outline — my original structure, with the three sub-topics kept
#  under "Perfect Competition".  Each row is (label, title, description,
#  is_sub).  The label is what goes in the gold circle.
#
#  No coverage pills and no video references in this round (2026-08-28,
#  Nico): the deck is built clean first, and the split into videos and
#  in-class material is a later pass.
# --------------------------------------------------------------------------

M4_OUTLINE = [
    ("1", "Introduction to market structures",
     "Where a firm sits between price taker and price setter", False),
    ("2", "Perfect competition",
     "Many small sellers, an identical product, and free entry", False),
    ("2a", "Profit maximization of a price taker in the short run",
     "How much to produce, and when to stop production", True),
    # 2026-08-29 (Nico): the old line ("Why the marginal cost curve is the
    # supply curve") was wrong — MC is the supply curve of ONE firm, not of
    # the market, which is the horizontal sum across firms.
    ("2b", "Firm-level and market supply",
     "Marginal cost is one firm’s supply curve; market supply sums across "
     "firms", True),
    ("2c", "Long-run competitive equilibrium",
     "Entry and exit drive economic profits to zero", True),
    ("3", "Market distortions and regulations",
     "Who wins, who loses, and what is lost outright", False),
    ("4", "Externalities",
     "Costs and benefits that land on people outside the deal", False),
]

# Which video each agenda item is taught in, per the course calendar's
# Session-4 prep block (2026-08-30).  A parent item that spans several
# videos names the range; its sub-items each name their single video.
# Items 4 and 5 are taught on campus in Part II.
COVERAGE_LABEL = {
    0: "Video 1",
    1: "Video 2",
    2: "Video 3",
    3: "Video 4",
    4: "Video 5",
    5: "In Class",
    6: "In Class",
}

# The video each item's CONTENT sits in, used for the top-bar tags.  None
# means the slides are taught in class and keep a two-level tag.
ITEM_VIDEO = {0: 1, 1: 2, 2: 3, 3: 4, 4: 5, 5: None, 6: None}

TAG_BASE = {i: _title_case(row[1][0].upper() + row[1][1:])
            for i, row in enumerate(M4_OUTLINE)}


def _tag_for(i):
    """Module 4 · Video k · <topic> inside a video block, two levels
    outside one (Teaching CLAUDE.md, top-bar tag rule)."""
    v = ITEM_VIDEO[i]
    return ("Module 4 \u00b7 Video %d \u00b7 %s" % (v, TAG_BASE[i])) if v \
        else ("Module 4 \u00b7 " + TAG_BASE[i])


TAG = {i: _tag_for(i) for i in range(len(M4_OUTLINE))}

TAG_INTRO = TAG[0]      # Introduction to Market Structures
TAG_PC = TAG[1]         # Perfect Competition
TAG_SR = TAG[2]         # Profit Maximization of a Price Taker in the Short Run
TAG_SUPPLY = TAG[3]     # Firm-Level and Market Supply
TAG_LR = TAG[4]         # Long-Run Competitive Equilibrium
TAG_DIST = TAG[5]       # Market Distortions and Regulations
TAG_EXT = TAG[6]        # Externalities

# how far a shaded (one-line) row is nudged down inside its reserved
# two-row box, so single-line rows sit centred — see Teaching CLAUDE.md
DIM_DROP = Inches(0.19)

# the coverage pill sits vertically centred in the reserved two-row
# box, so the pill column is identical on every agenda slide
PILL_DROP = Inches(0.16)


# ==========================================================================
#  Outline / agenda slides
# ==========================================================================

def make_m4_outline(prs, page_num, *, tag=None, title="Outline of Module 4",
                    descriptions=False, highlight_idx=None,
                    highlight_set=None, ps_pointer=False):
    """Module outline in the format converged on for Modules 1 – 3: a gold
    circle carrying the item label, the item title beside it in bold navy,
    and a one-line grey description underneath.

    Every item RESERVES the description row, so item positions are
    pixel-identical on every agenda slide.  The description shows only for
    the current topic(s), or for all of them when ``descriptions=True``.
    Section agendas band the current item in cream and shade the rest.

    Module 4 adds sub-items ("2a", "2b", "2c"): they are indented, set one
    step smaller, and their circle is a touch narrower, so the hierarchy of
    my original outline survives.
    """
    slide = _blank_slide(prs)
    # 2026-08-30: a section agenda carries ITS OWN section's video number
    # ("Module 4 · Video k · Agenda"); the descriptive overview belongs to
    # the introduction video, and an in-class section keeps two levels.
    if tag is None:
        v = ITEM_VIDEO.get(highlight_idx) if highlight_idx is not None else 1
        tag = ("Module 4 · Video %d · Agenda" % v) if v             else "Module 4 · Agenda"
    _draw_top_bar_tc(slide, tag)
    _draw_action_title(slide, title)

    hi = set()
    if highlight_idx is not None:
        hi.add(highlight_idx)
    if highlight_set:
        hi.update(highlight_set)
    if descriptions:
        hi = set(range(len(M4_OUTLINE)))

    n_items = len(M4_OUTLINE)
    top = Inches(1.42)
    bottom = Inches(7.02)
    gap = Inches(0.11) if n_items <= 6 else Inches(0.07)
    PITCH_MAX = Inches(0.91)          # the Module 2 pitch
    LAST_ROW_MAX = Inches(6.22)

    def _row_y(pitch, i):
        block = pitch * n_items - gap
        y0 = top + max(0, (bottom - top - block) // 2)
        return y0 + i * pitch

    pitch = PITCH_MAX
    while pitch > Inches(0.60) and (
            _row_y(pitch, n_items - 1) > LAST_ROW_MAX
            or pitch * n_items - gap > bottom - top):
        pitch -= 4572                  # 0.005" steps
    content = pitch - gap
    title_h = int(content * 0.525)
    desc_h = content - title_h
    y = int(_row_y(pitch, 0))

    for i, (label, item, desc, is_sub) in enumerate(M4_OUTLINE):
        lit = descriptions or i in hi
        circ_d = Inches(0.50) if is_sub else Inches(0.58)
        circ_x = Inches(1.60) if is_sub else Inches(1.15)
        text_x = Inches(2.35) if is_sub else Inches(2.05)
        t_size = 22 if is_sub else 25
        d_size = 20 if is_sub else 22

        if not descriptions and i in hi:
            band = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, int(Inches(0.90)),
                int(y - Inches(0.06)), int(Inches(12.15)),
                int(title_h + desc_h + Inches(0.10)))
            try:
                band.adjustments[0] = 0.35
            except Exception:
                pass
            band.fill.solid()
            band.fill.fore_color.rgb = CREAM
            band.line.color.rgb = GOLD
            band.line.width = Pt(1.0)
            band.shadow.inherit = False
            _add_drop_shadow(band)

        circ = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, int(circ_x), int(y + Inches(0.02)),
            int(circ_d), int(circ_d))
        circ.fill.solid()
        circ.fill.fore_color.rgb = GOLD
        circ.line.fill.background()
        circ.shadow.inherit = False
        tf = circ.text_frame
        tf.margin_left = tf.margin_right = 0
        tf.margin_top = tf.margin_bottom = 0
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = label
        run.font.size = Pt(19 if len(label) > 1 else 25)
        run.font.bold = True
        run.font.color.rgb = NAVY if lit else DIM
        run.font.name = "Calibri"

        rows = [([(_title_case(item[0].upper() + item[1:]),
                   {'bold': True, 'size': t_size,
                    'color': NAVY if lit else DIM})], 0,
                 {'bullet_style': 'none', 'space_before_pts': 0})]
        if i in hi:
            rows.append(([(desc, {'size': d_size, 'color': GRAY})], 0,
                         {'bullet_style': 'none', 'space_before_pts': 0}))
        # the coverage pill: gold = on video, navy = in class, and it
        # dims with its own row on a section agenda (Teaching CLAUDE.md).
        # A sub-item's pill is narrower but shares the right edge, so it
        # reads as part of its parent's.
        pill_w = Inches(1.14) if is_sub else Inches(1.55)
        pill = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, int(Inches(12.85) - pill_w),
            int(y + Inches(0.02) + PILL_DROP), int(pill_w),
            int(Inches(0.36)))
        try:
            pill.adjustments[0] = 0.30
        except Exception:
            pass
        label_txt = COVERAGE_LABEL[i]
        on_video = label_txt.startswith("Video")
        pill.fill.solid()
        pill.line.fill.background()
        pill.shadow.inherit = False
        if lit:
            pill.fill.fore_color.rgb = GOLD if on_video else NAVY
            pill_fg = NAVY if on_video else WHITE
            _add_drop_shadow(pill)
        else:
            pill.fill.fore_color.rgb = DIM
            pill_fg = WHITE
        ptf = pill.text_frame
        ptf.margin_left = ptf.margin_right = 0
        ptf.margin_top = ptf.margin_bottom = 0
        ptf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ppara = ptf.paragraphs[0]
        ppara.alignment = PP_ALIGN.CENTER
        prun = ppara.add_run()
        prun.text = label_txt
        prun.font.size = Pt(13)
        prun.font.bold = True
        prun.font.name = "Calibri"
        prun.font.color.rgb = pill_fg

        # the description must clear the pill by at least a five-letter
        # word, or it reads as running into it (course rule)
        if i in hi:
            avail = (Inches(12.85) - pill_w - Inches(0.10)) - text_x
            if _text_w_in(desc, d_size) > avail / 914400.0:
                raise ValueError(
                    "outline description %r is %.2f\" wide but only %.2f\" "
                    "clears the pill - shorten the wording"
                    % (desc, _text_w_in(desc, d_size), avail / 914400.0))

        _add_hierarchical_bullets(
            slide, text_x, y if lit else int(y + DIM_DROP),
            int(Inches(12.85) - pill_w - Inches(0.18)) - text_x,
            title_h + desc_h,
            rows, size=t_size, line_spacing_pts=0)
        y = int(y + pitch)

    _draw_footer(slide, FOOTER_TEXT, page_num)
    if ps_pointer:
        _add_ps_pointer(slide, top=Inches(6.68), label="Problem Set 3")
    return slide


# ==========================================================================
#  FRONT MATTER — slides 1 – 4
# ==========================================================================

def _video_title_card(prs, name, k):
    """A video title card: the title-slide layout with no top bar, no
    footer text and no page number (Teaching CLAUDE.md, "Converting a
    Module for Video Taping").  Carries no speaker notes."""
    slide = _blank_slide(prs)
    _add_text(slide, MARGIN, Inches(2.10), RULE_W, Inches(1.00), name,
              size=60, bold=True, color=NAVY, font="Calibri",
              align=PP_ALIGN.CENTER)
    _add_text(slide, MARGIN, Inches(3.25), RULE_W, Inches(0.70),
              "Module 4  \u00b7  Video %d" % k, size=40, bold=True,
              color=GOLD, font="Calibri", align=PP_ALIGN.CENTER)
    _add_rect(slide, (SLIDE_W - Inches(4.0)) // 2, Inches(4.28),
              Inches(4.0), Inches(0.05), GOLD)
    _add_text(slide, MARGIN, Inches(4.62), RULE_W, Inches(0.45),
              "Management 405", size=26, bold=True, color=GRAY,
              font="Calibri", align=PP_ALIGN.CENTER)
    _add_text(slide, MARGIN, Inches(5.32), RULE_W, Inches(0.40),
              "Prof. Nico Voigtl\u00e4nder  \u00b7  UCLA Anderson", size=22,
              color=GRAY, font="Calibri", align=PP_ALIGN.CENTER)
    _add_rect(slide, 0, Inches(7.15), SLIDE_W, Inches(0.02), RULE)
    _add_rect(slide, MARGIN, Inches(7.135), GOLD_W, Inches(0.05), GOLD)
    return slide


def slide_01_title(prs):
    """Deck title slide: centred, no top bar, no page number."""
    slide = _blank_slide(prs)
    # the deck name runs to two lines, so the box is tall enough for both
    _add_text(slide, MARGIN, Inches(1.95), RULE_W, Inches(1.70),
              DECK_TITLE, size=50, bold=True, color=NAVY,
              font="Calibri", align=PP_ALIGN.CENTER)
    _add_text(slide, MARGIN, Inches(3.72), RULE_W, Inches(0.80),
              DECK_SUB, size=40, bold=True, color=GOLD,
              font="Calibri", align=PP_ALIGN.CENTER)
    _add_rect(slide, (SLIDE_W - Inches(4.0)) // 2, Inches(4.62),
              Inches(4.0), Inches(0.06), GOLD)
    _add_text(slide, MARGIN, Inches(4.98), RULE_W, Inches(0.45),
              "Management 405", size=26, bold=True, color=GRAY,
              font="Calibri", align=PP_ALIGN.CENTER)
    _add_text(slide, MARGIN, Inches(5.52), RULE_W, Inches(0.45),
              "Prof. Nico Voigtländer  ·  UCLA Anderson", size=22,
              color=GRAY, font="Calibri", align=PP_ALIGN.CENTER)
    _add_rect(slide, MARGIN, Inches(7.15), RULE_W, Inches(0.02), RULE)
    _add_rect(slide, MARGIN, Inches(7.135), GOLD_W, Inches(0.05), GOLD)
    # 2026-08-29 (Nico): the UCLA Anderson wordmark that the original deck
    # carried in the lower-left corner (NV_s01_4_67f56784.jpg, 2.20" wide at
    # 0.60/6.42) is dropped.  It is UCLA blue, so it fought the navy/gold
    # palette, and the affiliation is already in the byline above.  No
    # institutional branding mark goes on any slide of this deck.
    return slide


def slide_02_logistics(prs, page_num):
    """Logistics.  Exam dates refreshed from the course calendar
    (2026-08-28, Nico): the Fall 2026 exam period is 11 – 13 December, and
    the slide must make clear that there is ONE 3.5-hour window, the same
    for everyone, with the exact date announced later."""
    bullets = [
        ("Coffee & Econ office hours restart next week", 0),
        ("Final exam period: December 11 – 13, 2026", 0),
        ("One 3.5-hour window — the same window for everyone", 1),
        ("The exact date and time will be announced later", 1),
        ("Online, open book, open notes; covers all material", 1),
    ]
    return content_slide(prs, page_num, TAG_LOG, "Some Logistics", bullets,
                         size=28, sub_size=24, line_spacing_pts=18)


def make_m4_roadmap(prs, page_num, *, tag=None):
    """Course roadmap in the Module-3 standard diamond format.  Module 4
    sits in box 4 ("Markets, Pricing, and Strategy"), so that box is navy
    and carries the gold 'we are here' arrow."""

    def draw(slide):
        box_h = Inches(0.85)
        narrow_w = Inches(4.6)
        wide_w = Inches(8.6)
        gap = Inches(0.3)
        mid = SLIDE_W // 2

        top_x = mid - wide_w // 2
        top_y = Inches(2.0)
        _add_rounded_filled_box(
            slide, top_x, top_y, wide_w, box_h,
            "1. Basic Principles and Economic Way of Thinking",
            fill=FADED, text_color=WHITE, size=24, bold=True)

        row2_y = Inches(3.65)
        left_x = mid - gap // 2 - narrow_w
        right_x = mid + gap // 2
        _add_rounded_filled_box(slide, left_x, row2_y, narrow_w, box_h,
                                "2. Value and Demand", fill=FADED,
                                text_color=WHITE, size=26, bold=True)
        _add_rounded_filled_box(slide, right_x, row2_y, narrow_w, box_h,
                                "3. Supply and Cost", fill=FADED,
                                text_color=WHITE, size=26, bold=True)

        bot_x = mid - wide_w // 2
        bot_y = Inches(5.5)
        _add_rounded_filled_box(slide, bot_x, bot_y, wide_w, box_h,
                                "4. Markets, Pricing, and Strategy",
                                fill=NAVY, text_color=WHITE, size=24,
                                bold=True)

        top_bottom_y = top_y + box_h
        _add_arrow(slide, (top_x + wide_w // 2, top_bottom_y),
                   (left_x + narrow_w // 2, row2_y),
                   color=FADED, weight_pt=3.0, head=True)
        _add_arrow(slide, (top_x + wide_w // 2, top_bottom_y),
                   (right_x + narrow_w // 2, row2_y),
                   color=FADED, weight_pt=3.0, head=True)
        row2_bottom_y = row2_y + box_h
        _add_arrow(slide, (left_x + narrow_w // 2, row2_bottom_y),
                   (bot_x + wide_w // 2, bot_y),
                   color=FADED, weight_pt=3.0, head=True)
        _add_arrow(slide, (right_x + narrow_w // 2, row2_bottom_y),
                   (bot_x + wide_w // 2, bot_y),
                   color=FADED, weight_pt=3.0, head=True)

        arrow_w, arrow_h = Inches(0.6), Inches(0.5)
        arrow_left = bot_x - arrow_w - Inches(0.12)
        arrow_top = bot_y + (box_h - arrow_h) // 2
        _add_arrow_shape(slide, arrow_left, arrow_top, arrow_w, arrow_h,
                         direction="right", fill=GOLD)
        _add_text(slide, arrow_left - Inches(1.55),
                  bot_y + (box_h - Inches(0.32)) // 2,
                  Inches(1.45), Inches(0.32), "we are here",
                  size=16, italic=True, bold=True, color=GOLD,
                  font="Calibri", align=PP_ALIGN.RIGHT)

    return make_diagram_slide(prs, page_num, tag or TAG_ROADMAP,
                              "Agenda for the Class", draw)


# ==========================================================================
#  1 · INTRODUCTION TO MARKET STRUCTURES — slides 5 – 7
# ==========================================================================

def slide_06_market_structures(prs, page_num):
    """The four basic market structures.  My original is six separate
    one-column tables butted together; this is ONE native table.  Monopoly
    products read "Unique" in quotation marks (2026-08-28, Nico)."""
    rows = [
        ["", "Perfect\nCompetition", "Monopolistic\nCompetition",
         "Oligopoly", "Monopoly"],
        ["Number of firms", "Many", "Many", "Few", "One"],
        ["Type of products sold", "Identical", "Differentiated",
         "Identical or\nDifferentiated", "“Unique”"],
        ["Barriers to entry", "None", "None", "Some", "Many"],
    ]
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_INTRO)
    _draw_action_title(slide, "The Four Basic Market Structures")
    tbl_w = Inches(12.2)
    first_col = Inches(3.0)
    other = (tbl_w - first_col) // 4
    _add_styled_table(
        slide, (SLIDE_W - tbl_w) // 2, Inches(2.10), tbl_w, Inches(3.40),
        rows, col_widths=[first_col] + [other] * 4,
        row_heights=[Inches(1.00)] + [Inches(0.80)] * 3,
        font_size=20, header_size=20)
    # 2026-08-29 (Nico): the takeaway reads left-to-right but the table runs
    # the other way — Monopoly on the right, Perfect Competition on the left
    # — so "more firms → more competition" is a move from RIGHT to LEFT.  The
    # bar is grown to hold a long leftward arrow underneath the sentence, so
    # the direction is on the slide rather than only in the telling.  (Was a
    # 9.60 x 0.55" bar at y 6.10 with no arrow.)
    bar_w, bar_h = Inches(11.60), Inches(1.02)
    bar_l = (SLIDE_W - bar_w) // 2
    bar_t = Inches(5.80)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, int(bar_l), int(bar_t),
        int(bar_w), int(bar_h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GOLD
    bar.line.fill.background()
    try:
        bar.adjustments[0] = 0.30
    except Exception:
        pass
    bar.shadow.inherit = False
    _add_drop_shadow(bar)
    tf = bar.text_frame
    tf.margin_left = tf.margin_right = Inches(0.10)
    tf.margin_top, tf.margin_bottom = Inches(0.05), Inches(0.12)
    tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "More firms, more similar products, easier entry → more competition"
    r.font.size = Pt(19)
    r.font.bold = True
    r.font.color.rgb = NAVY
    r.font.name = "Calibri"
    # 2026-08-29 (Nico): the arrow sits ABOVE the sentence inside the bar
    # (was below it), and is thicker — 5.0 pt, was 3.0.  It points back
    # towards the competitive (left) end of the table above.
    _add_arrow(slide, (bar_l + bar_w - Inches(0.55), Inches(6.12)),
               (bar_l + Inches(0.55), Inches(6.12)),
               color=NAVY, weight_pt=5.0, head=True, head_size='lg')
    _draw_footer(slide, FOOTER_TEXT, page_num)
    _set_notes(slide, NOTES[6])
    return slide


def slide_07_market_power(prs, page_num):
    """The price-taker → price-searcher spectrum, with one illustration per
    market structure and the 'TODAY' band over the left end."""
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_INTRO)
    _draw_action_title(
        slide, "Extent of Market Power: Price Takers vs. Price Searchers")

    # 2026-08-29 (Nico): reworked to adopt MW slide 8's treatment — the
    # "least / most market power" ends set large (28 pt bold navy, was 17 pt
    # italic grey), a DOUBLE-headed spectrum arrow, the price-taker /
    # price-searcher labels as filled pills, and dark-red column separators
    # and dark-red frame round Perfect Competition (MW uses C00000 for both).
    # The frame keeps my rounded corners but takes MW's 3.5 pt weight.  The
    # illustrations are blown up to fill their columns, the bottom takeaway
    # bar is dropped, and my TODAY pill stays in place of MW's "WE ARE HERE"
    # arrow.
    x0, x1 = Inches(0.50), Inches(12.83)
    cols = [Inches(0.55), Inches(3.60), Inches(6.70), Inches(9.75)]
    col_w = Inches(2.95)

    _add_text(slide, Inches(0.55), Inches(1.38), Inches(4.0), Inches(0.50),
              "Least Market Power", size=28, bold=True, color=NAVY,
              font="Calibri", align=PP_ALIGN.LEFT)
    _add_text(slide, Inches(8.78), Inches(1.38), Inches(4.0), Inches(0.50),
              "Most Market Power", size=28, bold=True, color=NAVY,
              font="Calibri", align=PP_ALIGN.RIGHT)
    _add_arrow(slide, (x0, Inches(2.02)), (x1, Inches(2.02)),
               color=NAVY, weight_pt=3.5, head_both=True, head_size='lg')

    labels = ["Perfect\nCompetition", "Monopolistic\nCompetition",
              "Oligopoly", "Monopoly"]
    for cx, lab in zip(cols, labels):
        _add_text(slide, cx, Inches(2.22), col_w, Inches(0.80), lab,
                  size=22, bold=True, color=NAVY, font="Calibri",
                  align=PP_ALIGN.CENTER)
    for sx in (Inches(3.45), Inches(6.55), Inches(9.60)):
        _add_rect(slide, sx, Inches(2.20), Inches(0.04), Inches(0.85), RED)

    _add_rounded_filled_box(slide, Inches(0.90), Inches(3.12),
                            Inches(2.25), Inches(0.44), "(Price Takers)",
                            fill=GOLD, text_color=NAVY, size=22,
                            corner_pct=0.35)
    _add_rounded_filled_box(slide, Inches(3.90), Inches(3.12),
                            Inches(8.65), Inches(0.44), "(Price Searchers)",
                            fill=GOLD, text_color=NAVY, size=22,
                            corner_pct=0.35)

    # the block we are in today: a dark-red outline round the left column,
    # sitting below the spectrum arrow so the arrow runs clear of it
    band = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, int(Inches(0.49)), int(Inches(2.14)),
        int(Inches(2.98)), int(Inches(4.65)))
    try:
        band.adjustments[0] = 0.04
    except Exception:
        pass
    band.fill.background()
    band.line.color.rgb = RED
    band.line.width = Pt(3.5)
    band.shadow.inherit = False

    # One illustration per structure, on a common baseline, each blown up to
    # fill 2.40" of column height (was 1.70").  The aircraft column stacks
    # two logos, so they share that height and are set to a COMMON WIDTH —
    # Boeing's 1.504 aspect against Airbus's 1.266 otherwise leaves Airbus
    # looking shrunken beside it.
    ph_top = Inches(3.72)
    ph_h = 2.40
    _add_media_image(slide, "NV_s07_10_f9f010a8.jpg",
                     left=Inches(0.78), top=ph_top, height=Inches(ph_h))
    _add_media_image(slide, "NV_s07_g1.1_1f7061fb.png",
                     left=Inches(3.79), top=ph_top, height=Inches(ph_h))
    jet_w = 1.594                       # Boeing at 1.06" tall; Airbus matches
    _add_media_image(slide, "NV_s07_g12.1_a0b89503.jpg",
                     left=Inches(8.175 - jet_w / 2), top=ph_top,
                     width=Inches(jet_w), rounded=False, shadow=False)
    _add_media_image(slide, "NV_s07_g12.2_3444df7e.png",
                     left=Inches(8.175 - jet_w / 2), top=Inches(4.86),
                     width=Inches(jet_w), rounded=False, shadow=False)
    _add_media_image(slide, "NV_s07_13_4dc4c7c4.jpg",
                     left=Inches(10.09), top=ph_top,
                     height=Inches(ph_h), rounded=False, shadow=False)

    captions = ["Commodity markets", "Restaurants in WeHo",
                "Wide-body aircraft", "Municipal utilities"]
    for cx, cap in zip(cols, captions):
        _add_text(slide, cx, Inches(6.20), col_w, Inches(0.32), cap,
                  size=13, italic=True, color=GRAY, font="Calibri",
                  align=PP_ALIGN.CENTER)

    # 2026-08-29 (Nico): back to straddling the MIDDLE of the frame's
    # BOTTOM edge, as it did before the slide-7 rework (t-1 had the pill
    # at 1.15/5.68 on a frame whose bottom edge was 5.90).  The frame
    # bottom is now 2.14 + 4.65 = 6.79", so the 0.46"-tall pill starts at
    # 6.56 and its centre x matches the frame's (0.44 + 2.98/2 = 1.93").
    # Fill is dark red at 30 % opacity so it reads as part of the frame.
    # An OPAQUE WHITE pill goes underneath first: without it the frame's
    # 3.5 pt line shows straight through the translucent fill and strikes
    # the word "TODAY" out.
    _add_rounded_filled_box(slide, Inches(1.13), Inches(6.56),
                            Inches(1.70), Inches(0.46), "",
                            fill=WHITE, text_color=WHITE, size=16,
                            corner_pct=0.30, shadow=False)
    _set_fill_alpha(
        _add_rounded_filled_box(slide, Inches(1.13), Inches(6.56),
                                Inches(1.70), Inches(0.46), "TODAY",
                                fill=RED, text_color=NAVY, size=16,
                                corner_pct=0.30), 30)
    _draw_footer(slide, FOOTER_TEXT, page_num)
    _set_notes(slide, NOTES[7])
    return slide


# ==========================================================================
#  2 · PERFECT COMPETITION — slides 8 – 11
# ==========================================================================

def slide_09_price_taker(prs, page_num):
    """Characteristics of a price-taker market.  Rebuilt as a native
    two-column table (adopted from MW slide 9) instead of two facing text
    columns."""
    rows = [
        ["", "Perfect competition — price takers"],
        ["Number of actual or\npotential competitors", "Many small sellers"],
        ["Product differentiation", "None"],
        ["Entry conditions", "No barriers"],
        ["Profit potential",
         "Short run: can be positive, negative or zero\n"
         "Long run: zero economic profit"],
        ["Examples",
         "Agricultural commodities; currency markets; unskilled labor"],
    ]
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_PC)
    _draw_action_title(slide, "Perfect Competition: The Price Taker")
    # 2026-08-29 (Nico): the table is narrowed from 11.80" (centred, cols
    # 4.30/7.50) to 8.40" pinned left, to free the right-hand third for the
    # wheat-farmers picture.  Column 2 keeps 5.50" so "Short run: can be
    # positive, negative or zero" (4.83" at 20 pt Calibri) still sits on one
    # line inside the 0.1" cell margins.
    tbl_w = Inches(8.40)
    _add_styled_table(
        slide, Inches(0.45), Inches(1.85), tbl_w, Inches(4.60),
        rows, col_widths=[Inches(2.90), Inches(5.50)],
        row_heights=[Inches(0.62), Inches(0.90), Inches(0.72),
                     Inches(0.72), Inches(1.00), Inches(0.80)],
        font_size=20, header_size=22)
    # square image, centred vertically on the table
    _add_media_image(slide, "Wheat Famers.png",
                     left=Inches(9.42), top=Inches(2.52),
                     width=Inches(3.40), height=Inches(3.40))
    _draw_footer(slide, FOOTER_TEXT, page_num)
    _set_notes(slide, NOTES[8])
    return slide


def slide_10_market_and_firm(prs, page_num):
    """Two panels: the market sets P*, and the individual firm faces that
    price as a horizontal demand curve d = MR."""

    def draw(slide):
        mk = SimpleFig(1.40, 6.35, 3.50, 3.55, xmax=10.0, ymax=10.0)
        fm = SimpleFig(7.90, 6.35, 3.80, 3.55, xmax=10.0, ymax=10.0)

        # panel titles sit above the y-axis titles, centred on their panel
        _add_text(slide, Inches(1.40), Inches(1.50), Inches(3.5),
                  Inches(0.40), "Market", size=22, bold=True, color=NAVY,
                  font="Calibri", align=PP_ALIGN.CENTER)
        _add_text(slide, Inches(7.90), Inches(1.50), Inches(3.8),
                  Inches(0.40), "Firm (price taker)", size=22, bold=True,
                  color=NAVY, font="Calibri", align=PP_ALIGN.CENTER)

        for fig in (mk, fm):
            _fig_axes(slide, fig, x_title="Quantity", y_title="Price ($)",
                      label_size=17)

        # market: S up, D down, equilibrium at (5, 5)
        # 2026-08-29 (Nico): demand is dark yellow and everything that
        # carries the market price over to the firm is dark red — the
        # equilibrium dot, the arrow pointing down onto it (adopted from MW
        # slide 10), the dashed carry-over line, the firm's d = MR line and
        # its label.  The short dashed segment from the P* axis label to the
        # equilibrium stays grey, so there is only ONE dark-red dashed line.
        _fig_line(slide, mk, (0.6, 1.0), (9.2, 9.0), color=NAVY,
                  weight_pt=2.5)
        _fig_line(slide, mk, (0.6, 9.0), (9.2, 1.0), color=DARKYELLOW,
                  weight_pt=2.5)
        _fig_curve_label(slide, mk, 9.35, 9.0, "S", size=20)
        _fig_curve_label(slide, mk, 9.35, 1.1, "D", size=20,
                         color=DARKYELLOW)
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, int(mk.x(4.9) - Inches(0.07)),
            int(mk.y(5.0) - Inches(0.07)), int(Inches(0.14)),
            int(Inches(0.14)))
        dot.fill.solid()
        dot.fill.fore_color.rgb = RED
        dot.line.color.rgb = RED
        dot.line.width = Pt(1.0)
        dot.shadow.inherit = False
        _add_arrow(slide, (mk.x(4.9), mk.y(5.0)), (fm.x(0.0), fm.y(5.0)),
                   color=RED, weight_pt=2.0, head=True, dash="dash")
        _fig_line(slide, mk, (0.0, 5.0), (4.9, 5.0), color=GRAY,
                  weight_pt=1.25, dash="dash")
        _fig_ylab(slide, mk, 5.0, "P*", size=18, bold=True)
        # label centred over the equilibrium, with an arrow dropping onto it
        _add_text(slide, mk.x(4.9) - Inches(0.95), Inches(2.68),
                  Inches(1.9), Inches(0.60), "Market\nequilibrium", size=15,
                  italic=True, color=RED, font="Calibri",
                  align=PP_ALIGN.CENTER)
        _add_arrow(slide, (mk.x(4.9), Inches(3.40)),
                   (mk.x(4.9), Inches(4.45)),
                   color=RED, weight_pt=2.0, head=True)

        # firm: horizontal d = MR at the market price
        _fig_line(slide, fm, (0.0, 5.0), (9.4, 5.0), color=RED,
                  weight_pt=3.0)
        # hand-tweaked 2026-08-29 (Nico): 4.42 -> 4.30, so the label
        # sits clear of the red d = MR line rather than on it
        _add_text(slide, Inches(6.87), Inches(4.30), Inches(0.95),
                  Inches(0.30), "P*", size=18, bold=True, italic=True,
                  color=NAVY, font="Calibri", align=PP_ALIGN.RIGHT)
        _add_text(slide, Inches(8.10), Inches(3.85), Inches(4.85),
                  Inches(0.40), "Demand curve  d  =  Marginal Revenue (MR)",
                  size=17, bold=True, color=RED, font="Calibri")

    slide = make_diagram_slide(prs, page_num, TAG_PC,
                               "Perfect Competition: The Market and the Firm",
                               draw)
    _set_notes(slide, NOTES[10])
    return slide


def slide_11_farmers(prs, page_num):
    """Germany 2024: farmers protest subsidy cuts because they cannot pass
    costs on.  The minister's line is the quote callout."""
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_PC)
    _draw_action_title(slide, "Why Farmers Cannot Pass On Their Costs")

    _add_hierarchical_bullets(
        slide, MARGIN, Inches(1.60), Inches(10.4), Inches(1.30),
        [("Germany cut its farm subsidies in 2024", 0),
         ("Some farmers stood to lose up to €10,000 a year", 0)],
        size=26, line_spacing_pts=14)

    # hand-tweaked in PowerPoint on 2026-08-29 (Nico): the photo group was
    # dragged out to 6.87 x 3.84" at (0.50, 2.93) and the quote box pulled in
    # to 5.09 x 2.45" at (7.76, 3.50).  Ported here as ungrouped shapes at
    # those coordinates.  The photo is kept at its native 1.689 aspect
    # (6.48 x 3.84" rather than his 6.87 x 3.84"), because the stretch was a
    # side effect of dragging a group corner, not a deliberate crop.
    # It also lands 6.44" rather than 6.87" wide so the caption underneath
    # clears the footer rule at y 7.15".
    # (Was: photo 5.60" wide at 0.50/3.05; caption at 0.50/6.48;
    #  quote box 6.05 x 2.35" at 6.80/3.60.)
    _add_media_image(slide, "NV_s11_4_e363d6f5.png",
                     left=Inches(0.50), top=Inches(2.90), width=Inches(6.44))
    _add_text(slide, Inches(0.50), Inches(6.78), Inches(6.44), Inches(0.32),
              "Farmer protests in Berlin, January 2024. "
              "Photo: The Guardian", size=12, italic=True, color=GRAY,
              font="Calibri", align=PP_ALIGN.CENTER)

    _add_convention_box(
        slide, Inches(7.76), Inches(3.50), Inches(5.09), Inches(2.45),
        runs=[("“For farmers there is a structural problem: they cannot "
               "pass on their production costs, because the prices are not "
               "made by them.”", {'italic': True, 'size': 19}),
              ("\n", {}),
              ("— Germany's Minister of Economic Affairs, 2024",
               {'bold': True, 'size': 17})],
        corner_pct=0.10, size=19)
    _add_media_image(slide, "NV_s11_5_d816cef7.png",
                     left=Inches(11.10), top=Inches(1.55), width=Inches(1.75),
                     rounded=False, shadow=False)
    _draw_footer(slide, FOOTER_TEXT, page_num)
    return slide


# ==========================================================================
#  Figure primitives added for Module 4
# ==========================================================================

def _fig_curve(slide, fig, fn, x0, x1, *, color=NAVY, weight_pt=2.5,
               segments=4, dash=None):
    """An editable freeform through fn over [x0, x1], built from a handful
    of cubic Bézier anchors (Catmull-Rom → Bézier), not a dense polyline,
    so "Edit Points" shows a few handles.  The shape gets its own TIGHT
    bounding box hugging just this curve, per the course chart rules.
    """
    n = segments
    pts = [(x0 + (x1 - x0) * i / n, None) for i in range(n + 1)]
    pts = [(xv, fn(xv)) for xv, _ in pts]
    dev = [(int(fig.x(xv)), int(fig.y(yv))) for xv, yv in pts]

    # Catmull-Rom control points
    ext = [dev[0]] + dev + [dev[-1]]
    beziers = []
    for i in range(1, len(ext) - 2):
        p0, p1, p2, p3 = ext[i - 1], ext[i], ext[i + 1], ext[i + 2]
        c1 = (p1[0] + (p2[0] - p0[0]) / 6.0, p1[1] + (p2[1] - p0[1]) / 6.0)
        c2 = (p2[0] - (p3[0] - p1[0]) / 6.0, p2[1] - (p3[1] - p1[1]) / 6.0)
        beziers.append((c1, c2, p2))

    xs = [p[0] for p in dev] + [c[0] for b in beziers for c in b[:2]]
    ys = [p[1] for p in dev] + [c[1] for b in beziers for c in b[:2]]
    left, top = int(min(xs)), int(min(ys))
    w, h = max(int(max(xs) - left), 1), max(int(max(ys) - top), 1)

    def rel(p):
        return int(p[0] - left), int(p[1] - top)

    parts = ["<a:moveTo><a:pt x=\"%d\" y=\"%d\"/></a:moveTo>" % rel(dev[0])]
    for c1, c2, p in beziers:
        parts.append(
            "<a:cubicBezTo>"
            "<a:pt x=\"%d\" y=\"%d\"/><a:pt x=\"%d\" y=\"%d\"/>"
            "<a:pt x=\"%d\" y=\"%d\"/></a:cubicBezTo>"
            % (rel(c1) + rel(c2) + rel(p)))

    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    spPr = shp._element.spPr
    old = spPr.find(qn('a:prstGeom'))
    dash_xml = ('<a:prstDash val="%s"/>' % dash) if dash else ''
    geom = _H.ET.fromstring(
        '<a:custGeom xmlns:a="%s"><a:avLst/><a:gdLst/><a:ahLst/>'
        '<a:cxnLst/><a:rect l="0" t="0" r="r" b="b"/>'
        '<a:pathLst><a:path w="%d" h="%d">%s</a:path></a:pathLst>'
        '</a:custGeom>' % (_H.A_NS, w, h, "".join(parts)))
    old.addnext(geom)
    spPr.remove(old)
    shp.fill.background()
    shp.shadow.inherit = False
    # Go through python-pptx's line API, NOT a bare SubElement: inside
    # <a:spPr> the order is xfrm -> geometry -> fill -> ln -> effectLst,
    # and an <a:ln> appended after the effect list is ignored outright
    # (the curve then renders in the theme's thin blue default).
    shp.line.color.rgb = color
    shp.line.width = Pt(weight_pt)
    ln = spPr.find(qn('a:ln'))
    # inside <a:ln> the order is fill -> dash -> join
    if dash_xml:
        ln.append(_H.ET.fromstring(
            dash_xml.replace('<a:prstDash',
                             '<a:prstDash xmlns:a="%s"' % _H.A_NS)))
    ln.append(_H.ET.fromstring('<a:round xmlns:a="%s"/>' % _H.A_NS))
    return shp


def _xlab_n(slide, fig, xv, label, *, size=17, bold=True, w=None,
            color=NAVY):
    """A NARROW x-axis label, sized to its own text (2026-08-30, Nico) so
    two quantities less than an inch apart never overlap."""
    if w is None:
        w = _text_w_in(label, size, bold=bold, italic=True) + 0.08
    return _add_text(slide, fig.x(xv) - Inches(w / 2),
                     Inches(fig.b + 0.06), Inches(w), Inches(0.30), label,
                     size=size, bold=bold, italic=True, color=color,
                     font="Calibri", align=PP_ALIGN.CENTER)


def _fig_dot(slide, fig, xv, yv, *, d=Inches(0.15), fill=GOLD, line=NAVY):
    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, int(fig.x(xv) - d // 2), int(fig.y(yv) - d // 2),
        int(d), int(d))
    dot.fill.solid()
    dot.fill.fore_color.rgb = fill
    dot.line.color.rgb = line
    dot.line.width = Pt(1.0)
    dot.shadow.inherit = False
    return dot


def _fig_region(slide, fig, x_lo, x_hi, y_lo, y_hi, *, fill=CREAM,
                alpha=None, line=None):
    """A shaded rectangle in logical coordinates (profit / loss areas)."""
    left, top = int(fig.x(x_lo)), int(fig.y(y_hi))
    w = int(fig.x(x_hi)) - left
    h = int(fig.y(y_lo)) - top
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1.25)
    shp.shadow.inherit = False
    if alpha is not None:
        sf = shp._element.spPr.find(qn('a:solidFill'))
        clr = sf.find(qn('a:srgbClr'))
        clr.append(_H.ET.fromstring(
            '<a:alpha xmlns:a="%s" val="%d"/>' % (_H.A_NS, int(alpha))))
    return shp


def _yi_badge(slide, label="Yi Family Example", left=9.95, top=6.52):
    """The corner tab that marks a slide as part of the running Yi-family
    example (my original deck carries it as loose text).

    2026-08-29 (Nico): MW's shape adopted — a FILLED rounded rect rather
    than my white outlined box — but rounded further than MW's 0.35 (to a
    full 0.50 pill) and filled in dark yellow at 35 % opacity, so it reads
    as a wash rather than a solid slab.  Was: white fill, 1.5 pt gold
    border, corner 0.28.
    """
    return _set_fill_alpha(
        _add_rounded_filled_box(
            slide, Inches(left), Inches(top), Inches(3.10), Inches(0.50),
            "🎬  " + label, fill=DARKYELLOW, text_color=NAVY,
            size=16, corner_pct=0.50, shadow=True), 35)


# ==========================================================================
#  2a · PROFIT MAXIMIZATION IN THE SHORT RUN — slides 12 – 41
# ==========================================================================

def slide_13_yi_family(prs, page_num):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_SR)
    _draw_action_title(slide, "Example: The Yi Family")
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.15), Inches(1.85), Inches(7.0),
        Inches(3.60),
        [("A South Korean family moves from California to rural Arkansas "
          "to start a farm", 0),
         ("Their goal: grow Korean produce and sell it to vendors in "
          "Dallas, at the market price", 0),
         ("We follow their cabbage crop — the raw material for kimchi", 0)],
        size=26, line_spacing_pts=20)
    # hand-tweaked 2026-08-29 (Nico): the whole right-hand picture column
    # moved up ~0.35" (poster 1.80 -> 1.45, caption 4.50 -> 4.15, the two
    # stills 4.95 -> 4.60 / 4.61)
    _add_media_image(slide, "NV_s12_3_534a7244.png",
                     left=Inches(8.05), top=Inches(1.45), width=Inches(4.60))
    _add_text(slide, Inches(8.05), Inches(4.15), Inches(4.60), Inches(0.30),
              "Minari (2020)", size=12, italic=True, color=GRAY,
              font="Calibri", align=PP_ALIGN.CENTER)
    _add_media_image(slide, "NV_s12_4_53752c12.jpg",
                     left=Inches(8.05), top=Inches(4.60), width=Inches(2.20))
    _add_media_image(slide, "NV_s12_5_042a9a61.png",
                     left=Inches(10.45), top=Inches(4.61), width=Inches(2.20))
    _yi_badge(slide)
    _draw_footer(slide, FOOTER_TEXT, page_num)
    _set_notes(slide, NOTES[12])
    return slide


def slide_14_cost_table(prs, page_num):
    """The cabbage cost schedule as a NATIVE table (my original is a
    spreadsheet screenshot), condensed to six rows — the row selection is
    adopted from MW slide 15, the numbers are mine."""
    qs = [50, 200, 400, 600, 800, 1000]
    rows = [["Q  (tons)", "TC", "TFC", "TVC"]]
    for q in qs:
        rows.append(["{:,}".format(q), "{:,.0f}".format(tc(q)),
                     "{:,.0f}".format(TFC),
                     "{:,.0f}".format(tc(q) - TFC)])
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_SR)
    _draw_action_title(slide, "Cabbage Production Costs")
    _add_styled_table(
        # 2026-08-29 (Nico): declared width was 6.30" while the four
        # columns sum to 6.60".  PowerPoint honours the COLUMNS, so the
        # table rendered 0.30" wider than the white backing card was
        # sized for and the shade sat off-register on the right edge.
        slide, Inches(0.75), Inches(1.95), Inches(6.60), Inches(4.20), rows,
        col_widths=[Inches(1.65)] * 4,
        row_heights=[Inches(0.66)] + [Inches(0.59)] * 6,
        font_size=19, header_size=19, first_col_align_left=False)

    _add_convention_box(
        slide, Inches(7.55), Inches(1.95), Inches(5.30), Inches(1.15),
        prefix="Fixed costs: ",
        body="land rent, equipment, and the salary the family gives up "
             "(opportunity cost)",
        corner_pct=0.12, size=17)
    _add_convention_box(
        slide, Inches(7.55), Inches(3.25), Inches(5.30), Inches(1.00),
        prefix="Variable costs: ",
        body="seeds, water, fertilizer, pesticides, fuel…",
        corner_pct=0.12, size=17)
    for fn, x in (("NV_s13_5_cd8d8b01.jpg", 7.55),
                  ("NV_s13_4_3f29b09c.jpg", 9.35),
                  ("NV_s13_1_dd9ab3ce.jpg", 11.15)):
        _add_media_image(slide, fn, left=Inches(x), top=Inches(4.50),
                         width=Inches(1.70), height=Inches(1.28))
    _add_text(slide, Inches(7.55), Inches(5.85), Inches(5.30), Inches(0.30),
              "Seeds  ·  irrigation  ·  fertilizer", size=12, italic=True,
              color=GRAY, font="Calibri", align=PP_ALIGN.CENTER)
    _yi_badge(slide)
    _draw_footer(slide, FOOTER_TEXT, page_num)
    _set_notes(slide, NOTES[13])
    return slide


def slide_15_tc_chart(prs, page_num):
    """Total cost against output: the observed cost data on the left, the
    fitted cost function through it on the right.

    2026-08-29 (Nico): reworked to MW slide 16's treatment — the six cost
    observations shown as DOTS with the estimated line running through them,
    and the equation boxed and labelled as the regression line.  The TFC
    guide line and its caption are dropped; TFC has its own column on the
    previous slide.

    NOTE on the data: MW's slide-16 chart plots about a dozen dots that
    scatter AROUND the curve at Q = 50, 100, 200 … 1100, and they do NOT
    reproduce the six-row table on MW slide 15 (that table's values sit
    exactly on TC = 60,000 + 135Q + 0.2Q²).  Rather than transcribe
    invented scatter off a picture, the dots here ARE the table on the
    left, so the two halves of the slide agree with each other.
    """
    qs = [50, 200, 400, 600, 800, 1000]

    def draw(slide):
        rows = [["Q  (tons)", "TC  ($)"]]
        for q in qs:
            rows.append(["{:,}".format(q), "{:,.0f}".format(tc(q))])
        _add_styled_table(
            slide, Inches(0.60), Inches(2.15), Inches(3.30), Inches(3.67),
            rows, col_widths=[Inches(1.50), Inches(1.80)],
            row_heights=[Inches(0.55)] + [Inches(0.52)] * 6,
            font_size=18, header_size=18, first_col_align_left=False)

        # xmax runs to 1,200 (not 1,100) so the "1,000" tick label
        # clears the "Quantity (tons)" axis title in the corner
        fig = SimpleFig(5.55, 6.05, 6.20, 3.35, xmax=1200.0, ymax=450000.0)
        _fig_axes(slide, fig, x_title="Quantity (tons)",
                  y_title="Total cost ($)", label_size=18)
        for xv in (200, 400, 600, 800, 1000):
            _fig_xlab(slide, fig, xv, "{:,}".format(xv), size=16)
        for yv in (100000, 200000, 300000, 400000):
            _fig_ylab(slide, fig, yv, "{:,}".format(yv), size=16)
        # the fitted line first, the observations on top of it
        _fig_curve(slide, fig, tc, 0, 1050, color=NAVY, weight_pt=3.0,
                   segments=4, dash="dash")
        for q in qs:
            _fig_dot(slide, fig, q, tc(q), d=Inches(0.17))

        # "Regression line:" over the equation, in one shape so the OMML is
        # never split off into a group of its own (see the grouping rule)
        rbox = _add_mixed_textbox(
            # Sits ABOVE the plot rather than inside it, as MW's does:
            # the fitted curve climbs through the plot's upper-left, so
            # a box in there is capped near 4.7" and OMML sets this
            # equation ~1.5x wider than the bare glyphs measure — it
            # would have had to drop below the 18 pt box-text floor.
            # Up here it clears the y-axis title (which ends at x 6.80)
            # and the equation gets its full 22 pt.
            slide, Inches(6.90), Inches(1.48), Inches(5.90), Inches(0.95),
            [("text", "Regression line:", {'size': 18, 'bold': True,
                                           'color': NAVY}),
             ("break", None, {}),
             ("omml", _omml_tc(), {'size': 22})],
            align=PP_ALIGN.CENTER)
        rbox.fill.solid()
        rbox.fill.fore_color.rgb = CREAM
        rbox.line.color.rgb = GOLD
        rbox.line.width = Pt(1.75)
        _apply_picture_style(rbox, corner_pct=10)

    slide = make_diagram_slide(prs, page_num, TAG_SR,
                               "Cabbage Total Cost Function", draw)
    _yi_badge(slide)
    _set_notes(slide, NOTES[14])
    return slide


def slide_16_business_relevance(prs, page_num):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_SR)
    _draw_action_title(slide, "Business Relevance")
    _add_text(slide, MARGIN + Inches(0.15), Inches(1.75), Inches(8.0),
              Inches(0.45), "Key questions:", size=28, bold=True,
              color=NAVY, font="Calibri")
    # 2026-08-29 (Nico): the two key questions adopt MW slide 12's format —
    # each sits in a cream rounded box with a 1.75 pt gold border, with the
    # numbered gold circle overlapping the box's left edge.  My wording and
    # the colourful "marginal analysis" line below are unchanged.  (Was: a
    # loose gold circle plus a plain 28 pt text line, no box.)
    for i, (num, q) in enumerate([
            ("1", "Should the Yi family produce at all?"),
            ("2", "How much should they produce?")]):
        box_y = Inches(2.42) + i * Inches(1.25)
        box = _add_rounded_filled_box(
            slide, Inches(1.05), box_y, Inches(11.30), Inches(0.95), q,
            fill=CREAM, text_color=NAVY, line=GOLD, size=26,
            corner_pct=0.12, shadow=True)
        box.line.width = Pt(1.75)
        circ = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, int(Inches(1.30)), int(box_y + Inches(0.20)),
            int(Inches(0.55)), int(Inches(0.55)))
        circ.fill.solid()
        circ.fill.fore_color.rgb = GOLD
        circ.line.fill.background()
        circ.shadow.inherit = False
        p = circ.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = num
        r.font.size = Pt(24)
        r.font.bold = True
        r.font.color.rgb = NAVY
        r.font.name = "Calibri"
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.15), Inches(4.85), Inches(12.4),
        Inches(1.30),
        [("We start with the second question: profit maximization", 0),
         ([("Key concept: ", {}),
           ("marginal analysis", {'bold': True, 'color': BLUE_PED})], 0, {})],
        size=28, line_spacing_pts=18)
    _yi_badge(slide)
    _draw_footer(slide, FOOTER_TEXT, page_num)
    _set_notes(slide, NOTES[15])
    return slide


def slide_17_profit_max_rule(prs, page_num):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_SR)
    _draw_action_title(slide, "Price Taker Profit Maximization")
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.15), Inches(1.80), Inches(12.4),
        Inches(3.10),
        [("Set MR = MC — the rule for every profit maximizer", 0),
         ("But MR is far easier to compute for a price taker", 0),
         ("For a price taker, MR = P, the market price", 1),
         ("P is given: the firm cannot set it", 1),
         ("So there is only one decision left: how much to produce", 0)],
        size=26, sub_size=24, line_spacing_pts=16)
    # 2026-08-29 (Nico): MW slide 14's format for the rule box — cream fill,
    # 1.75 pt GOLD border (was navy 0.75 pt), slight 0.12 rounding, and a
    # gold circle carrying MW's "➜" glyph overlapping the left edge.  The
    # wording and the native OMML setting of Q* are unchanged; MW's plain
    # text and its red "MC = P" are not adopted, so red keeps meaning
    # "the demand a single firm faces" across this deck.
    eq = _add_math_equation(
        slide, Inches(3.05), Inches(5.20), Inches(7.20), Inches(1.00),
        _omml_text('Find ') + _omml_sup(_omml_run('Q'), _omml_text('*'))
        # hand-tweaked 2026-08-29 (Nico): the RULE half is dark red and
        # bold; "Find Q* such that" stays navy
        + _omml_text(' such that  ') + _omml_run('P', color=RED, bold=True)
        + _omml_text(' = ', color=RED, bold=True)
        + _omml_text('MC', color=RED, bold=True),
        size_pt=30, color=NAVY, fill=CREAM, line=GOLD, rounded=True,
        shadow=True, corner_pct=12000)
    eq.line.width = Pt(1.75)
    badge = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, int(Inches(3.30)), int(Inches(5.45)),
        int(Inches(0.50)), int(Inches(0.50)))
    badge.fill.solid()
    badge.fill.fore_color.rgb = GOLD
    badge.line.fill.background()
    badge.shadow.inherit = False
    bp = badge.text_frame.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    br = bp.add_run()
    br.text = "➜"
    br.font.size = Pt(20)
    br.font.bold = True
    br.font.color.rgb = NAVY
    br.font.name = "Calibri"
    _draw_footer(slide, FOOTER_TEXT, page_num)
    _set_notes(slide, NOTES[16])
    return slide


def slide_18_revenue_conditions(prs, page_num):
    """The price-taker revenue table beside the two things it implies: TR
    rises along a straight ray, and d = P = MR is horizontal."""

    def draw(slide):
        rows = [["P", "Q", "TR", "MR"],
                ["$3", "10", "$30", "$3"],
                ["3", "11", "33", "3"],
                ["3", "12", "36", "3"],
                ["3", "13", "39", "3"]]
        _add_styled_table(
            slide, Inches(0.75), Inches(2.30), Inches(4.40), Inches(2.85),
            rows, col_widths=[Inches(1.10)] * 4,
            row_heights=[Inches(0.62)] + [Inches(0.56)] * 4,
            font_size=19, header_size=19, first_col_align_left=False)

        # 2026-08-29 (Nico): each panel gets a header, and his hand moves are
        # ported — the y-axis titles were pulled in tight to their text
        # (TR ($) to 5.78/2.17 at 0.70" wide, P ($) to 9.58/2.24 at 0.54"),
        # and the two in-plot equations were dropped to make room for the
        # headers (TR = P·Q 6.55/2.05 -> 7.30/3.08; d = P = MR 10.05/4.13 ->
        # 10.02/4.34).  y_title is therefore drawn by hand, not by _fig_axes.
        _add_text(slide, Inches(5.80), Inches(1.72), Inches(3.60),
                  Inches(0.36), "Total Revenue", size=20, bold=True,
                  color=NAVY, font="Calibri", align=PP_ALIGN.CENTER)
        _add_text(slide, Inches(9.55), Inches(1.72), Inches(3.60),
                  Inches(0.36), "The Firm's Demand", size=20, bold=True,
                  color=NAVY, font="Calibri", align=PP_ALIGN.CENTER)

        tr = SimpleFig(6.30, 5.95, 2.60, 3.30, xmax=15.0, ymax=48.0)
        _fig_axes(slide, tr, x_title="Quantity", y_title=None,
                  label_size=16)
        _add_text(slide, Inches(5.78), Inches(2.17), Inches(0.70),
                  Inches(0.27), "TR ($)", size=16, bold=True, italic=True,
                  color=NAVY, font="Calibri")
        _fig_line(slide, tr, (0, 0), (14.0, 42.0), color=NAVY,
                  weight_pt=2.75)
        _add_text(slide, Inches(7.30), Inches(3.08), Inches(2.6),
                  Inches(0.34), "TR  =  P · Q", size=18, bold=True,
                  italic=True, color=NAVY, font="Calibri")
        _fig_guide(slide, tr, (10, 30))
        _fig_xlab(slide, tr, 10, "10", size=15)
        _fig_ylab(slide, tr, 30, "30", size=15)

        # 2026-08-29 (Nico): the firm's demand line sits LOWER in the panel
        # — ymax 6.0 put P = 3 at exactly half the plot height, ymax 8.0
        # drops it to 37.5 % — and the line, its "3" tick and its label are
        # dark red, the colour this deck uses for the demand a single firm
        # faces (slide 10's d = MR line).  The label drops with the line.
        mr = SimpleFig(10.05, 5.95, 2.60, 3.30, xmax=15.0, ymax=8.0)
        _fig_axes(slide, mr, x_title="Quantity", y_title=None,
                  label_size=16)
        _add_text(slide, Inches(9.58), Inches(2.24), Inches(0.54),
                  Inches(0.27), "P ($)", size=16, bold=True, italic=True,
                  color=NAVY, font="Calibri")
        _fig_line(slide, mr, (0, 3), (14.0, 3), color=RED, weight_pt=3.0)
        _fig_ylab(slide, mr, 3, "3", size=15, color=RED)
        _add_text(slide, Inches(10.02), Inches(4.34), Inches(2.85),
                  Inches(0.34), "d  =  P  =  MR", size=18, bold=True,
                  italic=True, color=RED, font="Calibri",
                  align=PP_ALIGN.CENTER)

        _add_convention_box(
            slide, Inches(0.75), Inches(5.55), Inches(4.40), Inches(0.95),
            body="Selling one more unit always brings in the same $3, "
                 "so MR = P at every quantity",
            corner_pct=0.12, size=16)

    slide = make_diagram_slide(prs, page_num, TAG_SR,
                               "Price Taker Revenue Conditions", draw)
    return slide


def slide_19_tr_tc_visual(prs, page_num):
    """TR, TC and the profit curve together — rebuilt 2026-08-29 (Nico) to
    reproduce slide 18 of my ORIGINAL deck rather than the reduced two-curve
    version that stood here.

    Restored from the original: the green profit curve underneath, the two
    grey markers where the slopes match, the DRAWN tangent segments on TR
    and TC, "Slope of TR = slope of TC / MR = MC" with its arrow, the
    Profit = 0 crossing and the Maximum-profit callout, the Q* drop line,
    and the two small parenthetical notes.  Colours per his instruction:
    dark red TR, blue TC, dark green profit.  Axis labels are set a size
    smaller than the original's.  The one thing kept from my version is the
    cream "Profit = TR − TC" callout box.

    Geometry: TC = 2.2 + 0.30 Q + 0.1158 Q², TR = 1.6316 Q.  Chosen so the
    two cross at Q = 2 and Q = 9.5 (so profit is zero there and the profit
    parabola peaks between them) and the common-slope point is Q* = 5.75:
    1.6316 = 0.30 + 2(0.1158)Q*.  Every marker below is evaluated from
    these functions, never eyeballed.
    """
    # 2026-08-29 (Nico), second pass: a QUADRATIC TC cannot be made any more
    # convex than it already was.  For TC = F + aQ + bQ² the curve's bow away
    # from its own chord is b·Qmax²/4 against a total rise of b·Qmax², i.e.
    # exactly 1/4 whatever b is, and any a > 0 only flattens it — so raising b
    # again would just rescale the picture.  TC is therefore CUBIC now:
    #     TC = 5.0 + 0.10 Q + 0.030 Q³
    # whose bow is 0.385 of the rise rather than 0.25, and whose slope runs
    # from 0.10 at the origin to 9.1 at Q = 10.  The maximum TR - TC gap goes
    # from 3.80 to 7.96 (11.9 % -> 20.9 % of the plot height).
    #
    # The optimum moves with the curve and is solved, not assumed:
    # TR' = TC' gives p = a + 3cQ*², so Q* = sqrt((p - a) / 3c) = 6.0, and p
    # is set from that.  The first TR = TC crossing is no longer a round
    # number, so it is found by bisection rather than written down.
    TFC_V, A_LIN, C_CUB = 5.0, 0.10, 0.030
    QS = 6.0                                  # where the two slopes match
    P_SL = A_LIN + 3 * C_CUB * QS * QS        # = 3.34, the price / TR slope

    def tr_fn(q):
        return P_SL * q

    def tc_fn(q):
        return TFC_V + A_LIN * q + C_CUB * q ** 3

    def pi_fn(q):
        return tr_fn(q) - tc_fn(q)

    def _root(f, lo, hi, n=80):
        """Bisect f on [lo, hi], which must bracket a sign change."""
        for _ in range(n):
            mid = (lo + hi) / 2.0
            if f(lo) * f(mid) <= 0:
                hi = mid
            else:
                lo = mid
        return (lo + hi) / 2.0

    Q_ZERO = _root(pi_fn, 0.4, QS)             # first TR = TC crossing

    def draw(slide):
        fig = SimpleFig(2.55, 5.60, 6.30, 3.85, xmax=10.0, ymax=38.0)
        x0, ytop = fig.x(0), Inches(1.72)
        ybot = fig.y(-6.6)                    # the axis runs BELOW zero, as
        #                                       in the original, because the
        #                                       profit curve starts at -TFC
        _add_arrow(slide, (x0, ybot), (x0, ytop), color=NAVY, weight_pt=2.0,
                   head=True)
        _add_arrow(slide, (x0, fig.y(0)), (Inches(9.15), fig.y(0)),
                   color=NAVY, weight_pt=2.0, head=True)
        _add_text(slide, Inches(0.50), Inches(1.78), Inches(2.00),
                  Inches(1.10), "Firm revenue,\ncost and\nprofit ($)",
                  size=16, bold=True, italic=True, color=NAVY,
                  font="Calibri")
        _add_text(slide, Inches(9.22), Inches(5.62), Inches(1.70),
                  Inches(0.36), "Quantity (Q)", size=16, bold=True,
                  italic=True, color=NAVY, font="Calibri")

        # ---- the three curves ------------------------------------------
        _fig_line(slide, fig, (0, 0), (9.6, tr_fn(9.6)), color=RED,
                  weight_pt=3.0)
        _fig_curve(slide, fig, tc_fn, 0, 10.0, color=BLUE_PED,
                   weight_pt=3.5)
        _fig_curve(slide, fig, pi_fn, 0, 10.0, color=GREEN_DK,
                   weight_pt=3.0)
        _add_text(slide, Inches(8.05), Inches(1.72), Inches(0.90),
                  Inches(0.36), "TR", size=19, bold=True, italic=True,
                  color=RED, font="Calibri")
        _add_text(slide, Inches(8.88), Inches(2.05), Inches(0.90),
                  Inches(0.36), "TC", size=19, bold=True, italic=True,
                  color=BLUE_PED, font="Calibri")
        _add_text(slide, Inches(1.70), Inches(fig.y(TFC_V) / 914400 - 0.16),
                  Inches(0.80), Inches(0.32), "TFC", size=16, bold=True,
                  italic=True, color=BLUE_PED, font="Calibri",
                  align=PP_ALIGN.RIGHT)

        # ---- Q*: the drop line, the two markers, the tangent segments ---
        _fig_line(slide, fig, (QS, pi_fn(QS)), (QS, tr_fn(QS)), color=NAVY,
                  weight_pt=1.25, dash="dash")
        d_q = 0.80                            # half-length of a tangent bar
        for v in (tr_fn(QS), tc_fn(QS)):
            _fig_line(slide, fig, (QS - d_q, v - d_q * P_SL),
                      (QS + d_q, v + d_q * P_SL), color=NAVY, weight_pt=2.5)
        _fig_dot(slide, fig, QS, tr_fn(QS), d=Inches(0.17), fill=DIM,
                 line=NAVY)
        _fig_dot(slide, fig, QS, tc_fn(QS), d=Inches(0.17), fill=DIM,
                 line=NAVY)
        _fig_dot(slide, fig, QS, pi_fn(QS), d=Inches(0.17), fill=DIM,
                 line=NAVY)
        _add_text(slide, Inches(fig.x(QS) / 914400 - 0.45), Inches(5.66),
                  Inches(0.90), Inches(0.34), "Q*", size=18, bold=True,
                  italic=True, color=NAVY, font="Calibri",
                  align=PP_ALIGN.CENTER)

        # ---- the annotations, each with the arrow that points at it -----
        _add_text(slide, Inches(9.72), Inches(3.60), Inches(3.35),
                  Inches(0.74),
                  "Slope of TR = slope of TC\nMR = MC", size=15, bold=True,
                  italic=True, color=BLUE_PED, font="Calibri")
        _add_arrow(slide, (Inches(9.67), Inches(3.90)),
                   (fig.x(QS) + Inches(0.16), fig.y(tc_fn(QS))),
                   color=NAVY, weight_pt=2.0, head=True)

        _add_text(slide, Inches(9.72), Inches(4.60), Inches(2.60),
                  Inches(0.36), "Maximum profit", size=17, bold=True,
                  color=NAVY, font="Calibri")
        _add_arrow(slide, (Inches(9.67), Inches(4.78)),
                   (fig.x(QS) + Inches(0.16), fig.y(pi_fn(QS))),
                   color=NAVY, weight_pt=2.0, head=True)

        _fig_line(slide, fig, (Q_ZERO, 0), (Q_ZERO, tc_fn(Q_ZERO)),
                  color=NAVY, weight_pt=1.25, dash="dash")
        _add_text(slide, Inches(2.05), Inches(4.35), Inches(1.55),
                  Inches(0.36), "Profit = 0", size=17, bold=True,
                  color=NAVY, font="Calibri", align=PP_ALIGN.RIGHT)
        _add_arrow(slide, (Inches(3.30), Inches(4.72)),
                   (fig.x(Q_ZERO) - Inches(0.03),
                    fig.y(tc_fn(Q_ZERO)) - Inches(0.06)),
                   color=NAVY, weight_pt=2.0, head=True)

        _add_text(slide, Inches(3.30), Inches(2.52), Inches(2.30),
                  Inches(0.60),
                  "(TR has constant slope\nfor a price taker)", size=13,
                  italic=True, color=RED, font="Calibri")
        _add_text(slide, Inches(3.55), Inches(6.18), Inches(2.60),
                  Inches(0.60),
                  "(Profit is the difference\nbetween TR and TC)", size=13,
                  italic=True, color=GREEN_DK, font="Calibri")

        _add_convention_box(
            slide, Inches(9.72), Inches(1.80), Inches(3.33), Inches(1.45),
            runs=[("Profit = TR − TC", {'bold': True, 'size': 19}),
                  ("\n", {}),
                  ("Widest where the two slopes are equal "
                   "(MR = MC)", {'size': 16})],
            corner_pct=0.12, size=18)

    slide = make_diagram_slide(
        prs, page_num, TAG_SR,
        "Profit Maximization of a Price Taker: Visual Representation", draw)
    return slide


def slide_20_max_profit_setup(prs, page_num):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_SR)
    _draw_action_title(slide, "Maximizing Profit: Cabbage Production")
    # 2026-08-29 (Nico): the market price gets the same boxed treatment as
    # the cost function, so the two givens of the problem read as a pair
    # (it was the second of two plain bullets and disappeared).
    _add_text(slide, MARGIN + Inches(0.15), Inches(1.62), Inches(12.4),
              Inches(0.45), "The Yi family has estimated its total cost:",
              size=26, color=NAVY, font="Calibri")
    _add_math_equation(
        slide, Inches(3.55), Inches(2.20), Inches(6.2), Inches(0.95),
        _omml_tc(),
        size_pt=30, color=NAVY, fill=CREAM, line=NAVY, rounded=True,
        shadow=True)
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.15), Inches(3.28), Inches(12.4),
        Inches(0.45),
        [([("Q", {'italic': True}),
           (" is the quantity produced, in tons", {})], 0, {})],
        size=26, line_spacing_pts=0)
    _add_text(slide, MARGIN + Inches(0.15), Inches(3.86), Inches(12.4),
              Inches(0.45), "The market price of cabbage is:",
              size=26, color=NAVY, font="Calibri")
    _add_math_equation(
        slide, Inches(3.55), Inches(4.42), Inches(6.2), Inches(0.85),
        _omml_run('P') + _omml_text(' = $%s per ton' % _num(P_HIGH)),
        size_pt=30, color=NAVY, fill=CREAM, line=NAVY, rounded=True,
        shadow=True)
    # 2026-08-29 (Nico): the closing question comes out of the bullet list and
    # into a box of its own — pale blue fill (the deck's pedagogical blue at
    # 15 % opacity) with a 1.75 pt dark-blue border.
    qbox = _set_fill_alpha(
        _add_rounded_filled_box(
            slide, Inches(1.05), Inches(5.50), Inches(11.30), Inches(0.95),
            "How many tons should the Yi family produce to maximize profits?",
            fill=BLUE_PED, text_color=NAVY, line=BLUE_PED, size=26,
            corner_pct=0.12, shadow=True), 15)
    qbox.line.width = Pt(1.75)
    _yi_badge(slide)
    _draw_footer(slide, FOOTER_TEXT, page_num)
    _set_notes(slide, NOTES[19])
    return slide


def slide_22_qstar_solution(prs, page_num):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_SR)
    _draw_action_title(slide, "Profit-Maximizing Quantity: Solution")
    # 2026-08-29 (Nico): "MC is the derivative of TC with respect to Q" moves
    # to its OWN line (his hand edit), and two new steps are spelled out —
    # "MR = MC thus implies:" ahead of 135 + 0.4Q = 400, and "Solve for Q*"
    # ahead of the answer.  Eight items now share the content area, so
    # equations drop 28 -> 26 pt (the answer 32 -> 30).
    # 2026-08-29 (later): the full stop after "MR = MC" is dropped (his
    # hand edit) and the stack is re-pitched to read as FOUR STEPS rather
    # than eight evenly-spaced lines — each reason sits ~0.07" above the
    # equation it justifies, with ~0.22" between steps.
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.15), Inches(1.42), Inches(12.4),
        Inches(0.74),
        [([("We need to set ", {}), ("MR = MC", {'bold': True})], 0, {}),
         ([("MC is the derivative of TC with respect to ", {}),
           ("Q", {'italic': True}), (":", {})], 0, {})],
        size=24, line_spacing_pts=0)
    _add_math_equation(
        # 0.90" tall: the stacked dTC/dQ fraction is ~2 line-heights
        # and overflowed a 0.68" box up into the bullet above it
        slide, Inches(2.30), Inches(2.34), Inches(8.7), Inches(0.90),
        _omml_text('MC') + _omml_text(' = ')
        + _omml_frac(_omml_text('d') + _omml_text('TC'),
                     _omml_text('d') + _omml_run('Q'))
        + _omml_text(' = ') + _omml_mc_expr(),
        size_pt=26, color=NAVY)
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.15), Inches(3.32), Inches(12.4),
        Inches(0.38),
        [("The cabbage market is competitive, so MR is just the market "
          "price:", 0)],
        size=24, line_spacing_pts=0)
    _add_math_equation(
        slide, Inches(2.30), Inches(3.84), Inches(8.7), Inches(0.52),
        _omml_text('MR') + _omml_text(' = ') + _omml_run('P')
        + _omml_text(' = %s' % _num(P_HIGH)),
        size_pt=26, color=NAVY)
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.15), Inches(4.48), Inches(12.4),
        Inches(0.38),
        [([("MR = MC", {'bold': True}), (" thus implies:", {})], 0, {})],
        size=24, line_spacing_pts=0)
    _add_math_equation(
        slide, Inches(2.30), Inches(4.88), Inches(8.7), Inches(0.52),
        _omml_mc_expr() + _omml_text(' = %s' % _num(P_HIGH)),
        size_pt=26, color=NAVY)
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.15), Inches(5.64), Inches(12.4),
        Inches(0.38),
        [([("Solve for ", {}), ("Q", {'italic': True}), ("*", {}),
           ("  (the optimal production)", {})], 0, {})],
        size=24, line_spacing_pts=0)
    _add_math_equation(
        slide, Inches(2.32), Inches(6.16), Inches(8.7), Inches(0.58),
        _omml_sup(_omml_run('Q'), _omml_text('*'))
        + _omml_text(' = %s tons' % _num(Q_HIGH)),
        size_pt=30, color=DARKRED)
    _yi_badge(slide)
    _draw_footer(slide, FOOTER_TEXT, page_num)
    return slide


def slide_23_two_ways_atc(prs, page_num):
    """Adopted from MW slide 20: the profit identity written the two ways
    the deck goes on to use, stated before its first use."""
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_SR)
    # 2026-08-29 (Nico): retitled, and the gold "1" badge at 1.35/1.92
    # is dropped — it numbered a list that no longer exists
    _draw_action_title(slide, "Profits and Average Total Costs (ATC)")
    for i, eq in enumerate([
            _omml_text('Profit') + _omml_text(' = ') + _omml_text('TR')
            + _omml_text(' − ') + _omml_text('TC'),
            _omml_text('= (') + _omml_run('P') + _omml_text(' · ')
            + _omml_run('Q') + _omml_text(') − (') + _omml_text('ATC')
            + _omml_text(' · ') + _omml_run('Q') + _omml_text(')'),
            _omml_text('= (') + _omml_run('P') + _omml_text(' − ')
            + _omml_text('ATC') + _omml_text(') · ') + _omml_run('Q')]):
        _add_math_equation(
            slide, Inches(2.45), Inches(1.85) + i * Inches(0.92),
            Inches(8.4), Inches(0.80), eq, size_pt=30, color=NAVY)
    # the identity is only useful once ATC is on the table, so it is
    # defined here rather than assumed
    _add_math_equation(
        slide, Inches(2.45), Inches(4.72), Inches(8.4), Inches(1.00),
        _omml_text('where  ') + _omml_text('ATC') + _omml_text(' = ')
        + _omml_frac(_omml_text('TC'), _omml_run('Q'))
        + _omml_text('   (average total cost per ton)'),
        size_pt=24, color=GRAY)
    _add_math_equation(
        slide, Inches(2.45), Inches(5.90), Inches(8.4), Inches(0.85),
        _omml_text('Profit is positive if  ') + _omml_run('P')
        + _omml_text(' ≥ ') + _omml_text('ATC'),
        size_pt=28, color=NAVY, fill=CREAM, line=NAVY, rounded=True,
        shadow=True)
    _draw_footer(slide, FOOTER_TEXT, page_num)
    return slide


def slide_24_profit_solution(prs, page_num):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_SR)
    _draw_action_title(slide, "Profit at Q*: Solution")
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.15), Inches(1.68), Inches(12.4),
        Inches(0.80),
        [([("What is the economic profit at the profit-maximizing output "
            "level (", {}),
           ("Q* = %s" % _num(Q_HIGH), {'bold': True}),
           (") ?", {})], 0, {})],
        size=24, line_spacing_pts=0)
    eqs = [
        (_omml_text('Profit') + _omml_text(' = ') + _omml_text('TR')
         + _omml_text(' − ') + _omml_text('TC'), NAVY, 28),
        (_omml_text('TR') + _omml_text(' = ') + _omml_run('P')
         + _omml_text(' · ') + _omml_run('Q')
         + _omml_text(' = %s · %s = %s'
                      % (_num(P_HIGH), _num(Q_HIGH),
                         _num(P_HIGH * Q_HIGH))), NAVY, 26),
        (_omml_text('TC') + _omml_text(' = %s + %s · %s + %s · '
                                       % (_num(TFC), _num(B_LIN),
                                          _num(Q_HIGH), _num(B_QUAD)))
         + _omml_sup(_omml_text(_num(Q_HIGH)), _omml_text('2'))
         + _omml_text(' = %s' % _num(tc(Q_HIGH))), NAVY, 24),
        (_omml_text('Profit') + _omml_text(' = %s − %s = %s'
                                           % (_num(P_HIGH * Q_HIGH),
                                              _num(tc(Q_HIGH)),
                                              _num(P_HIGH * Q_HIGH
                                                   - tc(Q_HIGH)))),
         DARKRED, 30),
    ]
    for i, (eq, col, sz) in enumerate(eqs):
        _add_math_equation(
            slide, Inches(1.70), Inches(2.55) + i * Inches(0.88),
            Inches(10.0), Inches(0.80), eq, size_pt=sz, color=col)
    _add_takeaway_bar(
        slide, "The Yi family makes a positive economic profit at Q*",
        top=Inches(6.22), width=Inches(8.6), height=Inches(0.55),
        left=Inches(0.75), fill=GOLD, text_color=NAVY, size=19, bold=True,
        rounded=True, shadow=True)
    _yi_badge(slide)
    _draw_footer(slide, FOOTER_TEXT, page_num)
    return slide


def atc(q):
    return tc(q) / q


Q_ATC_MIN = (TFC / B_QUAD) ** 0.5          # 387.3 — where MC cuts ATC


def slide_25_profit_rectangle(prs, page_num):
    """The Yi family's profit as the rectangle between P and ATC.  Every
    marked point lies mathematically on the curves: MC and ATC cross at the
    minimum of ATC, and Q* is the true solution of MC = P."""

    def draw(slide):
        fig = SimpleFig(1.60, 6.00, 8.00, 3.85, xmax=950.0, ymax=500.0)
        _fig_axes(slide, fig, x_title="Quantity (tons)",
                  y_title="P", label_size=17)
        # The profit rectangle is drawn FIRST so the curves stay on top of
        # it — its lower edge runs along ATC, and ATC has to remain visible
        # through it.  Height P − ATC(Q*), width Q*.
        a_star = atc(Q_HIGH)
        # 2026-08-29 (Nico): the profit rectangle is dark red at 22 % opacity
        # (was solid cream), matching the callout on the right
        _fig_region(slide, fig, 0, Q_HIGH, a_star, P_HIGH,
                    fill=RED, alpha=22000, line=NAVY)

        _fig_line(slide, fig, (0, mc(0)), (800, mc(800)), color=NAVY,
                  weight_pt=2.75)
        _fig_curve_label(slide, fig, 778, mc(800) + 36, "MC", size=19)
        # MC hits the vertical axis at B_LIN: MC(0) = 135 + 2(0.2)(0) = 135
        _fig_ylab(slide, fig, mc(0), _num(mc(0)), size=15)
        # ATC starts low enough on its falling branch for the U to read;
        # any earlier and it would run off the top of the box
        _fig_curve(slide, fig, atc, 210, 800, color=GOLD, weight_pt=3.25,
                   segments=5)
        _fig_curve_label(slide, fig, 778, atc(800) - 44, "ATC", size=19,
                         color=GOLD)
        # the margin is genuinely thin (400 against 358), so the label sits
        # inside the band rather than the band being blown up to fit it
        _add_text(slide, Inches(2.30), fig.y(P_HIGH) + Inches(0.03),
                  Inches(2.4), Inches(0.32), "Profit", size=16, bold=True,
                  color=NAVY, font="Calibri", align=PP_ALIGN.CENTER)

        # the price line is the demand curve the firm faces, so it takes the
        # deck's dark red, label and axis value with it (2026-08-29, Nico)
        _fig_line(slide, fig, (0, P_HIGH), (860, P_HIGH), color=RED,
                  weight_pt=2.5)
        _fig_curve_label(slide, fig, 590, P_HIGH + 42, "P = MR", size=18,
                         color=RED)
        _fig_ylab(slide, fig, P_HIGH, _num(P_HIGH), size=17, bold=True,
                  color=RED)
        _fig_ylab(slide, fig, a_star, _num(round(a_star)), size=17)
        _fig_line(slide, fig, (Q_HIGH, 0), (Q_HIGH, P_HIGH), color=GRAY,
                  weight_pt=1.5, dash="dash")
        _fig_xlab(slide, fig, Q_HIGH, "Q* = %s" % _num(Q_HIGH),
                  size=17, bold=True)
        _fig_dot(slide, fig, Q_HIGH, P_HIGH)

        _add_mixed_textbox(
            slide, Inches(9.95), Inches(2.05), Inches(3.10), Inches(0.60),
            [("omml", _omml_text('MC') + _omml_text(' = ')
              + _omml_mc_expr(), {'size': 22})],
            align=PP_ALIGN.CENTER)

    slide = make_diagram_slide(
        prs, page_num, TAG_SR,
        "Illustrating Positive Profits of the Yi Family", draw)
    # the profit identity, with ATC written as a function so it is clear the
    # average cost is read off AT the optimum (2026-08-29, Nico)
    pbox = _add_mixed_textbox(
        slide, Inches(8.95), Inches(3.05), Inches(4.12), Inches(1.50),
        [("omml", _omml_text('Profit') + _omml_text(' = ')
          + _omml_text('(') + _omml_run('P') + _omml_text(' − ')
          + _omml_text('ATC') + _omml_text(') · ')
          + _omml_sup(_omml_run('Q'), _omml_text('*')), {'size': 18}),
         ("break", None, {}),
         ("text", "where ATC is evaluated at the optimal point Q*",
          {'size': 14, 'color': NAVY, 'font': 'Calibri'})],
        align=PP_ALIGN.CENTER)
    pbox.fill.solid()
    pbox.fill.fore_color.rgb = RED
    pbox.line.color.rgb = RED
    pbox.line.width = Pt(1.25)
    _apply_picture_style(pbox, corner_pct=12)
    _set_fill_alpha(pbox, 18)
    _yi_badge(slide)
    return slide


def slide_26_ross_stores(prs, page_num):
    """Where the cost concepts show up in a real income statement."""
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_SR)
    _draw_action_title(slide,
                       "Relationship to Accounting: Ross Stores Annual Report")
    # 2026-08-29 (Nico): rebuilt to follow my ORIGINAL slide 25 — the income
    # statement at full size with the cost concept written against each line,
    # rather than three small crops plus a separate glossary box.  The
    # statement is a table, so per the deck convention it is rebuilt NATIVELY
    # rather than shown as a screenshot; the figures are Ross Stores' as
    # filed, and both years foot exactly (13,708,907 + 2,874,469 + 74,328 =
    # 16,657,704 and 9,838,574 + 2,503,281 + 83,413 = 12,425,268).
    rows = [
        ["", "Economic cost concept", "Year ended\nJan 29, 2022",
         "Year ended\nJan 30, 2021"],
        ["Sales", "", "18,916,244", "12,531,565"],
        ["Cost of goods sold", "≈ TVC", "13,708,907", "9,838,574"],
        ["Selling, general and administrative", "Mix of FC & VC",
         "2,874,469", "2,503,281"],
        ["Interest expense (income), net", "Part of FC", "74,328", "83,413"],
        ["Total costs and expenses", "≈ TC", "16,657,704", "12,425,268"],
    ]
    tbl_w = Inches(12.00)
    _add_styled_table(
        slide, (SLIDE_W - tbl_w) // 2, Inches(1.95), tbl_w, Inches(3.75),
        rows, col_widths=[Inches(4.00), Inches(2.60), Inches(2.70),
                          Inches(2.70)],
        row_heights=[Inches(0.85)] + [Inches(0.58)] * 5,
        font_size=18, header_size=18,
        cell_text_colors={(r, 1): RED for r in range(1, 6)})
    _add_text(slide, (SLIDE_W - tbl_w) // 2, Inches(5.92), tbl_w,
              Inches(0.30),
              "Ross Stores, Form 10-K, fiscal 2021  ·  $ thousands",
              size=12, italic=True, color=GRAY, font="Calibri",
              align=PP_ALIGN.CENTER)
    _add_takeaway_bar(
        slide, "Accounting costs leave out opportunity costs",
        top=Inches(6.38), width=Inches(6.40), height=Inches(0.55),
        left=(SLIDE_W - Inches(6.40)) // 2, fill=GOLD, text_color=NAVY,
        size=18, bold=True, rounded=True, shadow=True)
    _draw_footer(slide, FOOTER_TEXT, page_num)
    _set_notes(slide, NOTES[25])
    return slide


# --------------------------------------------------------------------------
#  The general-case cost curves, shared by slides 27 and 35 – 37.
#
#  AVC is U-shaped, ATC sits above it, and MC cuts each at its minimum —
#  which the algebra below guarantees rather than approximates:
#      AVC(q) = a(q − q0)² + m      TC(q) = q·AVC(q) + F
#      MC(q)  = AVC(q) + q·AVC'(q) = a(3q² − 4q0·q + q0²) + m
#  so MC = AVC exactly at q = q0, and MC cuts ATC at ATC's minimum.
# --------------------------------------------------------------------------

#
#  The level constant _GC_M is the MINIMUM of AVC, and the shut-down case
#  needs the price line to sit visibly below it, so it is set well clear of
#  the axis rather than at a token value.
# 2026-08-29 (Nico): the U was too exaggerated.  Curvature drops
# 0.18 -> 0.12 and the drawing range widens (see GC_XLO / GC_XHI), so
# ATC is flatter but spans more of the box both across and up.
# 2026-08-30 (Nico): the general case now has the SHAPE of the Yi family's
# own cost function rather than the textbook U-shaped-AVC one.  With
#     TC = F + aQ + bQ²
# we get AVC = a + bQ and MC = a + 2bQ, both straight and RISING, while
# ATC = F/Q + a + bQ is the only U-shaped curve — and MC cuts ATC exactly at
# ATC's minimum, Q = sqrt(F/b), which the algebra below guarantees.
#
# ECONOMIC CONSEQUENCE, flagged rather than papered over: with a rising AVC,
# MC = a + 2bQ exceeds AVC = a + bQ at EVERY positive output, so a firm that
# sets P = MC always has P > AVC and never stops production.  The only way a
# firm with these costs produces nothing is P below the common intercept a,
# where the price line lies under AVC everywhere.  That is how the very-low-
# price panel is drawn.
_GC_A2, _GC_B2, _GC_F = 1.6, 0.55, 6.5

PINK = RGBColor(0xF3, 0xD8, 0xD8)        # legacy loss-region fill
GRAY_DK = RGBColor(0x3F, 0x44, 0x4C)     # the shut-down / stop-production
                                         # shading (see Teaching CLAUDE.md)


def gc_avc(q):
    return _GC_A2 + _GC_B2 * q


def gc_atc(q):
    return gc_avc(q) + _GC_F / q


def gc_mc(q):
    return _GC_A2 + 2 * _GC_B2 * q


def gc_q_star(p):
    """The output where MC = p.  Negative when p is below the intercept —
    the caller treats that as 'produce nothing'."""
    return (p - _GC_A2) / (2 * _GC_B2)


GC_XLO, GC_XHI = 1.20, 7.00      # ATC blows up near 0, so it starts at 1.2
GC_MC_HI = 7.00                  # MC is clipped so it stays inside the box
GC_AVC_HI = 7.00
GC_PRICE_HI = 7.30               # the price line stops here, leaving the
                                 # corner free for the d = MR label
GC_ATC_MIN_Q = (_GC_F / _GC_B2) ** 0.5      # = 3.438, where MC cuts ATC
GC_XMAX, GC_YMAX = 8.6, 9.5
GC_XMAX_SMALL = 9.4              # wider box on the three small panels, so
                                 # the Q* label clears the axis title
# HIGH: P > ATC(Q*) -> profit.  LOW: AVC(Q*) < P < ATC(Q*) -> loss but
# keep producing.  VLOW: P below the common MC / AVC intercept (1.6),
# the only case in which a firm with these costs produces nothing.
P_CASE_HIGH, P_CASE_LOW, P_CASE_VLOW = 7.5, 4.0, 1.15
P_CASE_NEG = 4.4                 # the loss panel of the general-case slide


def _draw_cost_panel(slide, fig, price, *, label, show_avc=True,
                     region_fill=None, region_alpha=None, x_title="Q",
                     price_color=None, d_mr=False):
    """MC / ATC (/ AVC) with a horizontal price line and the resulting Q*.

    The shaded profit / loss region is drawn FIRST, so the curves stay on
    top of it instead of being buried under the fill.

    ``price_color`` overrides the price line's colour (dark red where the
    line is being read as the firm's demand curve); ``d_mr=True`` adds the
    "d = MR" label at its right-hand end, per my original deck.
    """
    _fig_axes(slide, fig, x_title=x_title, y_title="P", label_size=17,
              titles_at_tip=True)
    qs = gc_q_star(price)
    p_col = price_color if price_color is not None else NAVY
    if region_fill is not None:
        lo, hi = sorted((price, gc_atc(qs)))
        _fig_region(slide, fig, 0, qs, lo, hi, fill=region_fill,
                    alpha=region_alpha, line=NAVY)
    _fig_curve(slide, fig, gc_mc, GC_XLO, GC_MC_HI, color=NAVY,
               weight_pt=3.0, segments=4)
    # 8 anchors, not 4: with only 4 the Bezier sags a little away
    # from the true ATC at Q*, and the profit / loss box - whose top
    # edge IS ATC(Q*) - then fails to meet the drawn curve
    _fig_curve(slide, fig, gc_atc, GC_XLO, GC_XHI, color=GOLD,
               weight_pt=3.0, segments=8)
    if show_avc:
        _fig_curve(slide, fig, gc_avc, GC_XLO, GC_AVC_HI, color=GRAY,
                   weight_pt=2.5, segments=4)
    _fig_line(slide, fig, (0, price), (GC_PRICE_HI, price), color=p_col,
              weight_pt=2.5)
    if d_mr:
        _add_text(slide, fig.x(GC_PRICE_HI + 0.15), fig.y(price) - Inches(0.17),
                  Inches(0.90), Inches(0.34), "d = MR", size=15, bold=True,
                  italic=True, color=p_col, font="Calibri")
    _fig_line(slide, fig, (qs, 0), (qs, price), color=GRAY,
              weight_pt=1.5, dash="dash")
    _fig_xlab(slide, fig, qs, "Q*", size=17, bold=True)
    _fig_ylab(slide, fig, price, label, size=17, bold=True, color=p_col)
    return qs


def slide_27_general_case(prs, page_num):
    """Positive, zero and negative profit in one row of three panels."""

    def draw(slide):
        # 2026-08-29 (Nico): one navy banner over the whole row, the three
        # headings pushed down under it, and the panels lowered to make the
        # room.  The price line is the demand curve the firm faces, so it is
        # dark red and labelled d = MR (from the original deck).  Profit is
        # shaded dark red and loss grey, each with its own π note written
        # INSIDE the shaded band rather than captioned below the panel.
        _add_rounded_filled_box(
            slide, Inches(0.55), Inches(1.46), Inches(12.54), Inches(0.52),
            "Profit Levels for Different Market Prices",
            fill=NAVY, text_color=WHITE, size=20, corner_pct=0.14)
        # the flatter ATC curve now cuts across the positive-profit band,
        # so that note is anchored at Q = 4.8 where the band is clear
        # ABOVE the curve; the loss band has no curve through it and
        # takes the plain midpoint.
        panels = [
            ("Positive Profit", P_CASE_HIGH, "π > 0", RED, RED,
             "P > ATC at Q*"),
            ("Zero Profit", gc_atc(GC_ATC_MIN_Q), None, None, NAVY,
             "P = ATC at Q*"),
            ("Negative Profit (Loss)", P_CASE_NEG, "π < 0", GRAY, NAVY,
             "P < ATC at Q*"),
        ]
        for i, (heading, price, note, fill, note_color,
                cond) in enumerate(panels):
            x = 0.55 + i * 4.32
            _add_rounded_filled_box(
                slide, Inches(x), Inches(2.10), Inches(3.90), Inches(0.52),
                heading, fill=NAVY, text_color=WHITE, size=18,
                corner_pct=0.14)
            fig = SimpleFig(x + 0.75, 6.10, 2.85, 2.81,
                            xmax=GC_XMAX_SMALL, ymax=GC_YMAX)
            qs = _draw_cost_panel(slide, fig, price, label="P",
                                  show_avc=False, region_fill=fill,
                                  region_alpha=22000, price_color=RED,
                                  d_mr=True)
            # hand-set 2026-08-30 (Nico): MC's label sits just off the top
            # of its curve, ATC's just past the end of its own
            _fig_curve_label(slide, fig, GC_MC_HI - 0.20,
                             gc_mc(GC_MC_HI) + 0.30, "MC", size=16)
            _fig_curve_label(slide, fig, GC_XHI + 0.22,
                             gc_atc(GC_XHI) - 0.15, "ATC", size=16,
                             color=GOLD)
            if note is not None:
                # hand-set 2026-08-30 (Nico): the note sits in the MIDDLE of
                # the shaded rectangle (not tucked against an edge), and the
                # profit one is dark red like the shading it labels
                lo, hi = sorted((price, gc_atc(qs)))
                _add_text(slide, fig.x(qs / 2) - Inches(0.45),
                          fig.y((lo + hi) / 2) - Inches(0.135),
                          Inches(0.90), Inches(0.27), note, size=16,
                          bold=True, color=note_color, font="Calibri",
                          align=PP_ALIGN.CENTER)
            # the condition that names the case, as a cream card in every
            # panel so all three read the same way
            _add_convention_box(
                slide, fig.x(0) + Inches(1.10), Inches(4.94),
                Inches(1.68), Inches(0.34), runs=[(cond, {'size': 15,
                                                          'bold': True})],
                corner_pct=0.20, size=15, border=GOLD,
                align=PP_ALIGN.CENTER, pad_h=Inches(0.04),
                pad_v=Inches(0.02))

    slide = make_diagram_slide(
        prs, page_num, TAG_SR,
        "Maximizing Profits in the Short Run: The General Case", draw)
    _add_takeaway_bar(
        slide, "Profit = (P − ATC) · Q*, so the sign of the profit is the "
               "sign of P − ATC",
        top=Inches(6.52), width=Inches(10.4), height=Inches(0.52),
        left=(SLIDE_W - Inches(10.4)) // 2, fill=GOLD, text_color=NAVY,
        size=18, bold=True, rounded=True, shadow=True)
    _set_notes(slide, NOTES[26])
    return slide


def slide_28_stop_producing(prs, page_num):
    """2026-08-29 (Nico): rebuilt from a bullet list into boxes, and the
    opening line ("Is a firm better off operating at a loss, or shutting
    down...") dropped as confusing.  The short-run identity that used to be
    block 2 of the following slide moves ONTO this slide, so the horizon
    distinction and the formula the shut-down rule rests on sit together;
    the old "Two Ways to Write Profit" slide is deleted.
    """
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_SR)
    # hand-edited title and banner 2026-08-30 (Nico), ported verbatim apart
    # from the "Prodution" typo and the mismatched closing quote
    _draw_action_title(
        slide, "If Profit Is Negative, Should a Firm Stop Production / "
               "Shut Down?")

    _add_rounded_filled_box(
        slide, Inches(0.70), Inches(1.50), Inches(11.93), Inches(0.62),
        "“Shutting down” in the long run is different from "
        "“Stopping production” in the short run",
        fill=NAVY, text_color=WHITE, size=22, corner_pct=0.12)

    # 2026-08-30 (Nico): SHORT run on the left (it is the case we go on to
    # analyse), long run on the right.  Both boxes are the same height and
    # tall enough for the longer of the two texts.
    # 2026-08-30 (Nico): his wording, plus a blank line under each heading
    # and TOP anchoring so "Short Run" and "Long Run" sit at the same height
    # however long the two bodies are.  The fixed-cost line in each box is
    # dark red — it is the hinge the whole distinction turns on.
    # hand-set 2026-08-30 (Nico), read through the group transform this
    # time: the LEFT card renders 5.967" wide at x 0.700 and the RIGHT one
    # 5.590" wide at x 7.040 — he resized the GROUPS, which is why reading
    # the child shapes alone kept showing the original 5.80.
    for x, w, runs in (
            (0.700, 5.967, [("Short Run", {'bold': True, 'size': 24}),
                    ("\n", {'bold': True, 'size': 19}),
                    ("Temporarily halt production, with the option to "
                     "restart later\n", {'size': 17}),
                    ("Fixed costs still have to be paid\n",
                     {'size': 17, 'bold': True, 'color': RED}),
                    ("— e.g. a building rent or a business truck lease",
                     {'size': 17})]),
            (7.040, 5.590, [("Long Run", {'bold': True, 'size': 24}),
                    ("\n", {'bold': True, 'size': 19}),
                    ("Exit the industry altogether, shut down\n",
                     {'size': 17}),
                    ("No further payment of fixed costs\n",
                     {'size': 17, 'bold': True, 'color': RED}),
                    ("— e.g. sell the truck and let the building lease "
                     "expire", {'size': 17})])):
        # hand-tweaked 2026-08-30 (Nico): pad_h 0 so the inner text box is
        # the FULL 5.80" of the card (it was inset 0.20" each side, which
        # squeezed the wrapping), pad_v 0.12 as before
        _add_convention_box(slide, Inches(x), Inches(2.20), Inches(w),
                            Inches(1.95), runs=runs, corner_pct=0.12,
                            size=17, align=PP_ALIGN.CENTER,
                            anchor=MSO_ANCHOR.TOP, pad_h=0,
                            space_before_pts=6)

    _add_takeaway_bar(
        slide, "We analyze the short run’s stopping-production "
               "decision first:",
        top=Inches(4.26), width=Inches(11.93), height=Inches(0.58),
        left=Inches(0.70), fill=GOLD, text_color=NAVY,
        size=20, bold=True, rounded=True, shadow=True)

    # the short-run form of profit, with the two cost terms braced together
    _add_math_equation(
        slide, Inches(1.20), Inches(4.92), Inches(10.1), Inches(1.00),
        _omml_text('Short run:   ') + _omml_text('Profit') + _omml_text(' = ')
        + _omml_run('P') + _omml_text(' · ') + _omml_run('Q')
        + _omml_text(' − ')
        + _omml_underbrace(
            _omml_text('TVC') + _omml_text(' − ') + _omml_text('TFC'),
            _omml_text('TVC') + _omml_text(' + ') + _omml_text('TFC')
            + _omml_text(' = ') + _omml_text('TC')),
        size_pt=26, color=NAVY)
    # the working form: P in dark red, and AVC / TFC each in their own
    # colour so the arrow below can pick TFC out (2026-08-30, Nico)
    _add_math_equation(
        slide, Inches(1.30), Inches(6.07), Inches(6.60), Inches(0.57),
        _omml_text('Profit:   ') + _omml_run('π') + _omml_text(' = (')
        + _omml_run('P', color=RED) + _omml_text(' − ')
        + _omml_text('AVC', color=DARKYELLOW) + _omml_text(') · ')
        + _omml_run('Q')
        + _omml_text(' − ') + _omml_text('TFC', color=RED),
        size_pt=24, color=NAVY, fill=WHITE, line=NAVY, rounded=True,
        shadow=True, corner_pct=12000)
    # hand-tweaked 2026-08-30 (Nico): longer, and turned round so it points
    # AT the TFC term rather than away from it
    _add_arrow(slide, (Inches(8.85), Inches(6.59)),
               (Inches(7.29), Inches(6.36)), color=RED, weight_pt=2.5,
               head=True)
    _set_fill_alpha(
        _add_rounded_filled_box(
            slide, Inches(8.85), Inches(6.46), Inches(3.60), Inches(0.50),
            "Fixed costs still have to be paid", fill=RED,
            text_color=NAVY, line=RED, size=17, bold=True,
            corner_pct=0.14,
            shadow=False), 15)
    _draw_footer(slide, FOOTER_TEXT, page_num)
    return slide


def slide_30_shutdown_rule(prs, page_num):
    """My slide 28, with the operate / shut-down rule set as a gold bar
    (the framing is adopted from MW slide 26)."""
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_SR)
    _draw_action_title(
        slide, "The Short Run Decision to Stop Production")
    # 2026-08-30 (Nico): the "Compare what the firm earns..." lead-in is
    # deleted and everything moves up into the room it leaves.
    # 2026-08-29 (Nico): the premise of the whole comparison, stated
    # before the two options rather than left implicit.  Everything
    # below is re-pitched to make room.
    # 2026-08-30 (Nico): short-run wording throughout — the firm "produces"
    # or "stops production"; "shut down" is reserved for the long run.
    _add_convention_box(
        slide, Inches(1.45), Inches(1.58), Inches(10.45), Inches(0.60),
        runs=[("In the short run, fixed costs have to be paid even if the "
               "firm does not produce", {'bold': True, 'size': 20})],
        corner_pct=0.14, size=20, align=PP_ALIGN.CENTER)
    # Option 1 is the STATUS QUO (keep producing) and option 2 the change,
    # so that option 1 − option 2 is (P − AVC) · Q with the sign the rule
    # below is stated in.
    _add_rounded_filled_box(
        slide, Inches(0.70), Inches(2.40), Inches(5.85), Inches(0.55),
        "Option 1: continue to produce", fill=NAVY, text_color=WHITE,
        size=20, corner_pct=0.12)
    _add_rounded_filled_box(
        slide, Inches(6.90), Inches(2.40), Inches(5.85), Inches(0.55),
        "Option 2: stop production", fill=NAVY, text_color=WHITE, size=20,
        corner_pct=0.12)
    _add_math_equation(
        slide, Inches(0.70), Inches(3.29), Inches(5.85), Inches(0.50),
        _omml_text('Profit:   ') + _omml_run('π')
        + _omml_text(' = (') + _omml_run('P')
        + _omml_text(' − ') + _omml_text('AVC') + _omml_text(') · ')
        + _omml_run('Q') + _omml_text(' − ') + _omml_text('TFC'),
        size_pt=24, color=NAVY)
    _add_math_equation(
        slide, Inches(6.90), Inches(3.28), Inches(5.85), Inches(0.54),
        _omml_text('Profit:  ') + _omml_run('π')
        + _omml_text(' = − ') + _omml_text('TFC'),
        size_pt=24, color=NAVY)
    _add_text(slide, Inches(0.70), Inches(4.02), Inches(5.85), Inches(0.42),
              "Producing adds (P − AVC) · Q on top", size=19, italic=True,
              color=GRAY, font="Calibri", align=PP_ALIGN.CENTER)
    _add_text(slide, Inches(6.90), Inches(4.02), Inches(5.85), Inches(0.42),
              "The fixed costs are lost either way", size=19, italic=True,
              color=GRAY, font="Calibri", align=PP_ALIGN.CENTER)
    # both options feed the same comparison, so an arrow comes in from each
    _add_arrow(slide, (Inches(3.62), Inches(4.55)),
               (Inches(5.35), Inches(5.02)), color=NAVY, weight_pt=2.25,
               head=True)
    _add_arrow(slide, (Inches(9.82), Inches(4.55)),
               (Inches(8.09), Inches(5.02)), color=NAVY, weight_pt=2.25,
               head=True)
    # his notation: capital pi for profit, so the gap between the two
    # options is a change in profit, delta-Pi
    _add_mixed_textbox(
        slide, Inches(1.60), Inches(5.10), Inches(10.1), Inches(0.77),
        [("text", "Difference in profits when producing:",
          {'size': 22, 'color': NAVY, 'font': 'Calibri'}),
         ("break", None, {}),
         ("omml", _omml_text('\u0394') + _omml_run('π')
          + _omml_text(' = ') + _omml_run('Q') + _omml_text(' · (')
          + _omml_run('P') + _omml_text(' − ') + _omml_text('AVC')
          + _omml_text(')'), {'size': 26})],
        align=PP_ALIGN.CENTER)
    _set_fill_alpha(
        _add_rounded_filled_box(
            slide, (SLIDE_W - Inches(10.2)) // 2, Inches(6.30),
            Inches(10.2), Inches(0.62),
            "\u21d2   Rule: Continue to produce if  P ≥ AVC        ·        "
            "Stop production if  P < AVC",
            fill=RED, text_color=NAVY, line=RED, size=22,
            corner_pct=0.30, shadow=False), 22)
    _draw_footer(slide, FOOTER_TEXT, page_num)
    return slide


def slide_31_new_price(prs, page_num):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_SR)
    _draw_action_title(slide,
                       "Optimizing in the Short Run: A New Market Price")
    # 2026-08-30 (Nico): name the OLD price too, set the new one apart,
    # drop the box round TC, and use MW slide 27's question format.
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.15), Inches(1.60), Inches(12.4),
        Inches(1.30),
        [("The US starts importing large quantities of cabbage from China", 0),
         ([("The US market price drops from $%s to " % _num(P_HIGH), {}),
           ("$%s per ton" % _num(P_LOW),
            {'bold': True, 'size': 28, 'color': RED})], 0, {})],
        size=25, line_spacing_pts=14)
    _add_mixed_textbox(
        slide, MARGIN + Inches(0.15), Inches(3.05), Inches(12.4),
        Inches(0.55),
        [("text", "Remember:   ", {'size': 24, 'color': NAVY}),
         ("omml", _omml_tc(), {'size': 24})])
    _add_text(slide, MARGIN + Inches(0.15), Inches(3.86), Inches(12.4),
              Inches(0.42), "Questions", size=25, bold=True,
              color=NAVY, font="Calibri")
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.55), Inches(4.42), Inches(11.2),
        Inches(1.90),
        [("1.   What quantity should the Yi family now produce?", 0,
          {'bullet_style': 'none'}),
         ("2.   What profit do they make at that quantity?", 0,
          {'bullet_style': 'none'}),
         ("3.   Should they continue to produce in the short run?", 0,
          {'bullet_style': 'none'})],
        size=24, line_spacing_pts=16)
    _yi_badge(slide, left=10.12, top=6.10)
    _draw_footer(slide, FOOTER_TEXT, page_num)
    # the Poll Break badge is drawn LAST so it sits in front of the footer
    # rule it straddles (Teaching CLAUDE.md, 2026-08-30)
    _add_pollbreak_badge(slide)
    _set_notes(slide, NOTES[29])
    return slide


def slide_33_new_qstar(prs, page_num):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_SR)
    _draw_action_title(slide, "New Q* and Profit: Solution")
    # 2026-08-30 (Nico): "Solution 1" labels the FIRST line — the whole
    # MC = MR = P step is the answer to question 1 — and "Solution 2"
    # starts where the profit computation starts, at TR.  Profit is written
    # with the pi symbol throughout.
    eqs = [
        (_omml_text('MC') + _omml_text(' = ')
         + _omml_text('MR') + _omml_text(' = ') + _omml_run('P'), NAVY, 26),
        (_omml_mc_expr()
         + _omml_text(' = %s' % _num(P_LOW))
         + _omml_text('     →     ')
         + _omml_sup(_omml_run('Q', color=RED), _omml_text('*',
                                                           color=RED))
         + _omml_text(' = %s' % _num(Q_LOW), color=RED), NAVY, 26),
        (_omml_text('TR')
         + _omml_text(' = ') + _omml_run('P')
         + _omml_text(' · ') + _omml_run('Q')
         + _omml_text(' = %s · %s = %s'
                      % (_num(P_LOW), _num(Q_LOW),
                         _num(P_LOW * Q_LOW))), NAVY, 26),
        (_omml_text('TC') + _omml_text(' = %s + %s · %s + %s · '
                                       % (_num(TFC), _num(B_LIN),
                                          _num(Q_LOW), _num(B_QUAD)))
         + _omml_sup(_omml_text(_num(Q_LOW)), _omml_text('2'))
         + _omml_text(' = %s' % _num(tc(Q_LOW))), NAVY, 24),
        (_omml_run('π')
         + _omml_text(' = %s − %s = −%s'
                      % (_num(P_LOW * Q_LOW), _num(tc(Q_LOW)),
                         _num(tc(Q_LOW) - P_LOW * Q_LOW))),
         DARKRED, 28),
    ]
    # 2026-08-30 (Nico): the five lines read as TWO blocks — question 1 and
    # its answer, then question 2 and its three working lines — so the pitch
    # inside a block is tight and the gap between them is wide.
    # 2026-08-30 (Nico): each block is headed by its OWN left-aligned
    # "Question n:" label, so the two answers are visually separated
    # instead of the label being buried in the first equation.
    for y, lab in ((1.38, "Question 1:"), (3.06, "Question 2:")):
        _add_text(slide, MARGIN + Inches(0.15), Inches(y), Inches(4.0),
                  Inches(0.38), lab, size=24, bold=True, color=NAVY,
                  font="Calibri")
    tops = (1.80, 2.38, 3.48, 4.02, 4.60)
    for (eq, col, sz), y in zip(eqs, tops):
        _add_math_equation(
            slide, Inches(1.55), Inches(y), Inches(10.3), Inches(0.70),
            eq, size_pt=sz, color=col)
    _add_takeaway_bar(
        slide, "At the new price, the Yi family makes a loss at the "
               "optimal output level",
        top=Inches(5.35), width=Inches(9.2), height=Inches(0.55),
        left=Inches(0.75), fill=GOLD, text_color=NAVY, size=19, bold=True,
        rounded=True, shadow=True)
    _yi_badge(slide, left=10.12, top=6.10)
    _draw_footer(slide, FOOTER_TEXT, page_num)
    _add_pollbreak_badge(slide)
    return slide


def slide_34_operate_solution(prs, page_num):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_SR)
    # 2026-08-30 (Nico): the SOLUTION is walked through first — AVC, its
    # value at Q*, the price, and the comparison — and only then, in a
    # separate box of its own, the direct option-by-option check.  The old
    # "P = 210 > AVC = 172.5" line stated the answer and the question at
    # once, so it is split into "recall the price" and "is P > AVC?".
    _draw_action_title(slide, "Decision to Continue Production: Solution")
    _add_math_equation(
        slide, Inches(1.55), Inches(1.56), Inches(10.3), Inches(1.00),
        _omml_text('AVC') + _omml_text(' = ')
        + _omml_frac(_omml_text(_num(B_LIN)) + _omml_run('Q')
                     + _omml_text(' + %s' % _num(B_QUAD))
                     + _omml_sup(_omml_run('Q'), _omml_text('2')),
                     _omml_run('Q'))
        + _omml_text(' = ') + _omml_avc_expr(),
        size_pt=26, color=NAVY)
    _add_math_equation(
        slide, Inches(1.55), Inches(2.66), Inches(10.3), Inches(0.78),
        _omml_text('AVC') + _omml_text(' at ')
        + _omml_sup(_omml_run('Q'), _omml_text('*'))
        + _omml_text(' = %s  is  %s + %s = %s'
                     % (_num(Q_LOW), _num(B_LIN), _num(B_QUAD * Q_LOW),
                        _num(avc(Q_LOW)))),
        size_pt=26, color=NAVY)
    _add_math_equation(
        slide, Inches(1.55), Inches(3.50), Inches(10.3), Inches(0.68),
        _omml_text('Recall:   ') + _omml_run('P')
        + _omml_text(' = %s' % _num(P_LOW)),
        size_pt=26, color=NAVY)
    _add_math_equation(
        slide, Inches(1.55), Inches(4.22), Inches(10.3), Inches(0.72),
        _omml_text('Is   ') + _omml_run('P') + _omml_text('  >  ')
        + _omml_text('AVC'),
        size_pt=28, color=DARKRED)
    # 2026-08-30 (Nico): the conclusion is the deck's blue box, centred,
    # with a double arrow; the alternative route below it is cream.
    _set_fill_alpha(
        _add_rounded_filled_box(
            slide, (SLIDE_W - Inches(7.60)) // 2, Inches(5.02),
            Inches(7.60), Inches(0.60),
            "\u21d2   Keep producing in the short run",
            fill=BLUE_PED, text_color=NAVY, line=BLUE_PED, size=21,
            corner_pct=0.14, shadow=False), 15)
    # hand-placed 2026-08-30 (Nico): pulled left and narrowed, and
    # "Alternatively" underlined
    altbox = _add_mixed_textbox(
        slide, Inches(0.28), Inches(5.91), Inches(6.24), Inches(1.10),
        [("text", "Alternatively",
          {'size': 22, 'color': NAVY, 'underline': True}),
         ("text", ", compare the two options directly:",
          {'size': 22, 'color': NAVY}),
         ("break", None, {}),
         ("text", "\u25aa   Continue to produce:  loss of %s"
                  % _num(tc(Q_LOW) - P_LOW * Q_LOW),
          {'size': 20, 'color': NAVY}),
         ("break", None, {}),
         ("text", "\u25aa   Stop production:  loss of %s, the whole of TFC"
                  % _num(TFC), {'size': 20, 'color': NAVY})])
    altbox.fill.solid()
    altbox.fill.fore_color.rgb = CREAM
    altbox.line.color.rgb = GOLD
    altbox.line.width = Pt(1.25)
    _apply_picture_style(altbox, corner_pct=10)
    _yi_badge(slide, left=10.12, top=6.10)
    _draw_footer(slide, FOOTER_TEXT, page_num)
    _add_pollbreak_badge(slide)
    return slide


# --------------------------------------------------------------------------
#  A CUBIC cost function, for the "more complex costs" example (2026-08-30).
#
#      TC = F + aQ - bQ^2 + cQ^3
#      AVC = a - bQ + cQ^2      (U-shaped, minimum at Q = b/2c)
#      MC  = a - 2bQ + 3cQ^2    (U-shaped, minimum at Q = b/3c, i.e. to the
#                                LEFT of AVC's, and cutting AVC at AVC's own
#                                minimum -- which the algebra guarantees)
#      ATC = F/Q + AVC          (U-shaped)
#
#  This is the shape that MAKES the stop-production case drawable: MC dips
#  below AVC over a visible range, so a price under AVC's minimum still has
#  a Q* where MC = P, and the loss there is bigger than TFC.  The quadratic
#  costs used on slides 34 - 36 cannot do that (MC > AVC at every Q > 0).
#  AVC's left branch is deliberately shallower than MW slide 32's, where AVC
#  sits almost on top of ATC.
# --------------------------------------------------------------------------
_CX_A, _CX_B, _CX_C, _CX_F = 8.0, 2.4, 0.28, 3.0
CX_XLO, CX_XHI = 1.15, 6.10
CX_XMAX, CX_YMAX = 7.6, 10.5
# AVC's minimum is 2.857, ATC's is 3.536.  HIGH sits above ATC's minimum
# (profit), MID between the two minima (loss, keep producing) and LOW below
# AVC's minimum (stop production).
CX_P_HIGH, CX_P_MID, CX_P_LOW = 5.2, 3.2, 2.2
CX_PRICE = CX_P_LOW


def cx_avc(q):
    return _CX_A - _CX_B * q + _CX_C * q * q


def cx_mc(q):
    return _CX_A - 2 * _CX_B * q + 3 * _CX_C * q * q


def cx_atc(q):
    return _CX_F / q + cx_avc(q)


def cx_q_star(p):
    """Where MC = p on MC's RISING branch."""
    a, b, c = 3 * _CX_C, -2 * _CX_B, _CX_A - p
    return (-b + (b * b - 4 * a * c) ** 0.5) / (2 * a)


def _complex_cost_slide(prs, page_num, title, price, *, kind, footnote):
    """One price case on the CUBIC cost function — the same three panels as
    slides 36 - 38, but with a U-shaped AVC so MC really does dip below it.
    Shading follows the deck convention: profit dark red, ordinary loss
    grey, stop-production dark grey."""
    fill, note = {"profit": (RED, "Profit"),
                  "loss": (GRAY, "Loss"),
                  "shutdown": (GRAY_DK, "Loss")}[kind]

    def draw(slide):
        _add_rounded_filled_box(
            slide, Inches(2.90), Inches(1.46), Inches(7.55), Inches(0.56),
            "Example with more complex cost functions",
            fill=NAVY, text_color=WHITE, size=20, corner_pct=0.14)
        fig = SimpleFig(3.05, 5.50, 7.35, 3.10,
                        xmax=CX_XMAX, ymax=CX_YMAX)
        _fig_axes(slide, fig, x_title="Q", y_title="P", label_size=17,
                  titles_at_tip=True)
        qs = cx_q_star(price)
        # the loss if the firm produced: width Q*, height ATC(Q*) - P
        lo, hi = sorted((price, cx_atc(qs)))
        _fig_region(slide, fig, 0, qs, lo, hi,
                    fill=fill, alpha=25000, line=NAVY)
        _fig_curve(slide, fig, cx_mc, CX_XLO, CX_XHI, color=NAVY,
                   weight_pt=3.0, segments=8)
        _fig_curve(slide, fig, cx_atc, CX_XLO, CX_XHI, color=GOLD,
                   weight_pt=3.0, segments=8)
        _fig_curve(slide, fig, cx_avc, CX_XLO, CX_XHI, color=GRAY,
                   weight_pt=2.5, segments=8)
        _fig_line(slide, fig, (0, price), (CX_XHI + 0.55, price),
                  color=RED, weight_pt=2.5)
        _add_text(slide, fig.x(CX_XHI + 0.70), fig.y(price)
                  - Inches(0.17), Inches(0.90), Inches(0.34), "d = MR",
                  size=15, bold=True, italic=True, color=RED,
                  font="Calibri")
        _fig_ylab(slide, fig, price, "P", size=17, bold=True, color=RED)
        _fig_line(slide, fig, (qs, 0), (qs, price), color=GRAY,
                  weight_pt=1.5, dash="dash")
        _fig_xlab(slide, fig, qs, "Q*", size=17, bold=True)
        _fig_curve_label(slide, fig, CX_XHI - 0.55, cx_mc(CX_XHI) + 0.55,
                         "MC", size=18)
        _fig_curve_label(slide, fig, CX_XHI + 0.05, cx_atc(CX_XHI) + 0.45,
                         "ATC", size=18, color=GOLD)
        _fig_curve_label(slide, fig, CX_XHI + 0.05, cx_avc(CX_XHI) - 0.55,
                         "AVC", size=18, color=GRAY)
        # hand-set 2026-08-30 (Nico) on slide 37: the word sits at the LEFT
        # end of the shaded rectangle rather than centred in it, and the
        # profit one is dark red like the shading it labels (navy for a
        # loss, as on slide 27).
        _add_text(slide, fig.x(0) + Inches(0.145),
                  fig.y((lo + hi) / 2) - Inches(0.152),
                  Inches(1.40), Inches(0.303), note, size=18, bold=True,
                  color=(DARKRED if kind == "profit" else NAVY),
                  font="Calibri", align=PP_ALIGN.CENTER)
        cbox = _add_convention_box(
            slide, Inches(1.00), Inches(6.10), Inches(11.30), Inches(0.80),
            runs=footnote,
            corner_pct=0.12, size=18, border=GOLD, align=PP_ALIGN.CENTER)
        cbox.line.width = Pt(1.75)

    return make_diagram_slide(prs, page_num, TAG_SR, title, draw)


def slide_37b_complex_profit(prs, page_num):
    return _complex_cost_slide(
        prs, page_num, "High Price: Positive Profits in the Short Run",
        CX_P_HIGH, kind="profit",
        footnote=[("Optimal output is where ", {}),
                  ("MC = MR = P", {'bold': True}),
                  (".  This occurs at Q*.  At that output positive "
                   "economic profits are being made because ", {}),
                  ("P > ATC", {'bold': True, 'color': RED}),
                  (".", {})])


def slide_37c_complex_loss(prs, page_num):
    return _complex_cost_slide(
        prs, page_num, "Low Price: Producing at a Loss in the Short Run",
        CX_P_MID, kind="loss",
        footnote=[("Negative economic profits (losses) are being made "
                   "because ", {}),
                  ("P < ATC", {'bold': True, 'color': RED}),
                  (".  But ", {}),
                  ("P > AVC", {'bold': True, 'color': RED}),
                  ("   →   the firm should continue to produce in the "
                   "short run.", {})])


def slide_37d_complex_shutdown(prs, page_num):
    return _complex_cost_slide(
        prs, page_num,
        "Very Low Price: Stopping Production in the Short Run",
        CX_P_LOW, kind="shutdown",
        footnote=[("MC really does run below AVC for a while, so there "
                   "IS a Q* where MC = P.  But ", {}),
                  ("P < AVC at Q*", {'bold': True, 'color': RED}),
                  ("   →   producing would lose more than TFC, so the "
                   "firm stops production.", {})])


def _poll_placeholder(prs, page_num, tag):
    """A bare PollEverywhere slide: the top-bar tag, the footer, and the
    round POLL pill bottom-right.  No action title and NO title rule — the
    live poll is pasted onto the empty canvas (2026-08-30, Nico)."""
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, tag)
    _draw_footer(slide, FOOTER_TEXT, page_num)
    _draw_poll_pill(slide, position='bottom-right', fill=GOLD,
                    text_color=NAVY, dot_color=None, width=Inches(1.49),
                    height=Inches(0.51), text_size=20, shadow=True)
    return slide


def _price_case_slide(prs, page_num, title, price, price_label, *,
                      region_kind, footnote, avc_dy=-0.40):
    """One of the three price cases on the general cost curves."""

    def draw(slide):
        # shading per the Teaching CLAUDE.md convention (2026-08-30, Nico):
        # profit dark red, ordinary loss grey, stop-production DARK grey
        fill, note = {
            "profit": (RED, "Profit"),
            "loss": (GRAY, "Loss"),
            "shutdown": (GRAY_DK, "Loss"),
        }[region_kind]
        # 2026-08-30 (Nico): the panel is centred on the slide now that the
        # right-hand column is gone, and its axis titles sit at the arrow
        # tips per the Teaching CLAUDE.md rule.
        fig = SimpleFig(3.05, 5.50, 7.35, 3.55,
                        xmax=GC_XMAX, ymax=GC_YMAX)
        if region_kind == "shutdown":
            # P sits below the common MC / AVC intercept, so MC = P has no
            # positive solution: the firm's best output is ZERO and its loss
            # is TFC.  There is therefore no rectangle of width Q* and
            # height |P - ATC(Q*)| to shade — the earlier version drew one
            # with a flat top at AVC(7.0), which measured nothing.  The
            # panel now says the thing that IS true: the price line lies
            # below AVC at every output.
            _fig_axes(slide, fig, x_title="Q", y_title="P",
                      label_size=17, titles_at_tip=True)
            _fig_curve(slide, fig, gc_mc, 0.0, GC_MC_HI, color=NAVY,
                       weight_pt=3.0, segments=4)
            _fig_curve(slide, fig, gc_atc, GC_XLO, GC_XHI, color=GOLD,
                       weight_pt=3.0, segments=4)
            _fig_curve(slide, fig, gc_avc, 0.0, GC_AVC_HI, color=GRAY,
                       weight_pt=2.5, segments=4)
            _fig_line(slide, fig, (0, price), (GC_PRICE_HI, price),
                      color=RED, weight_pt=2.5)
            _add_text(slide, fig.x(GC_PRICE_HI + 0.15),
                      fig.y(price) - Inches(0.17), Inches(0.90),
                      Inches(0.34), "d = MR", size=15, bold=True,
                      italic=True, color=RED, font="Calibri")
            _fig_ylab(slide, fig, price, price_label, size=17, bold=True,
                      color=RED)
        else:
            _draw_cost_panel(slide, fig, price, label=price_label,
                             region_fill=fill, region_alpha=22000,
                             price_color=RED, d_mr=True)
        _fig_curve_label(slide, fig, GC_MC_HI - 0.45,
                         gc_mc(GC_MC_HI) + 0.65, "MC", size=18)
        _fig_curve_label(slide, fig, GC_XHI + 0.20,
                         gc_atc(GC_XHI) + 0.40, "ATC", size=18, color=GOLD)
        _fig_curve_label(slide, fig, GC_AVC_HI + 0.20,
                         gc_avc(GC_AVC_HI) + avc_dy, "AVC", size=18,
                         color=GRAY)
        # the word goes INSIDE the shaded rectangle: its width is 0..Q* and
        # its height runs between the price line and ATC(Q*), so its AREA is
        # the profit or loss the panel is about
        if region_kind != "shutdown":
            qs = gc_q_star(price)
            lo, hi = sorted((price, gc_atc(qs)))
            _add_text(slide, fig.x(qs / 2) - Inches(0.70),
                      fig.y((lo + hi) / 2) - Inches(0.19),
                      Inches(1.40), Inches(0.38), note, size=18, bold=True,
                      color=NAVY, font="Calibri", align=PP_ALIGN.CENTER)
        else:
            _add_text(slide, fig.x(GC_XLO + 0.4),
                      fig.y(price) + Inches(0.10), Inches(3.60),
                      Inches(0.34), "P < AVC at every output", size=17,
                      bold=True, italic=True, color=GRAY_DK,
                      font="Calibri")
        # MW's format: one cream bar under the chart, the rule in bold navy
        # and each price comparison in bold dark red
        cbox = _add_convention_box(
            slide, Inches(1.00), Inches(6.10), Inches(11.30), Inches(0.80),
            runs=footnote, corner_pct=0.12, size=18, border=GOLD,
            align=PP_ALIGN.CENTER)
        cbox.line.width = Pt(1.75)

    slide = make_diagram_slide(prs, page_num, TAG_SR, title, draw)
    return slide


def slide_35_high_price(prs, page_num):
    slide = _price_case_slide(
        prs, page_num, "High Price: Positive Profits in the Short Run",
        P_CASE_HIGH, "P high", region_kind="profit",
        footnote=[("Optimal output is where ", {}),
                  ("MC = MR = P", {'bold': True}),
                  (".  This occurs at Q*.  At that output positive "
                   "economic profits are being made because ", {}),
                  ("P > ATC", {'bold': True, 'color': RED}),
                  (".", {})])
    _set_notes(slide, NOTES[33])
    return slide


def slide_36_low_price(prs, page_num):
    return _price_case_slide(
        prs, page_num, "Low Price: Producing at a Loss in the Short Run",
        P_CASE_LOW, "P low", region_kind="loss", avc_dy=-0.75,
        footnote=[("Negative economic profits (losses) are being made "
                   "because ", {}),
                  ("P < ATC", {'bold': True, 'color': RED}),
                  (".  But ", {}),
                  ("P > AVC", {'bold': True, 'color': RED}),
                  ("   →   the firm should continue to produce in the "
                   "short run.", {})])


def slide_37_very_low_price(prs, page_num):
    return _price_case_slide(
        prs, page_num, "Very Low Price: Stopping Production in the Short Run",
        P_CASE_VLOW, "P very low", region_kind="shutdown", avc_dy=-0.75,
        footnote=[("There is no output at which producing covers "
                   "variable cost: ", {}),
                  ("P < AVC everywhere", {'bold': True, 'color': RED}),
                  ("   →   the firm should stop production, and its loss "
                   "is then ", {}),
                  ("TFC", {'bold': True, 'color': RED}),
                  (".", {})])


def slide_38_coffee(prs, page_num):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_SR)
    _draw_action_title(slide, "Coffee Bean Producer")
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.15), Inches(1.75), Inches(7.6),
        Inches(1.30),
        [("Your firm is a price taker in the coffee bean market", 0),
         ([("You have just ", {}),
           ("optimized", {'bold': True, 'underline': True}),
           (" your output level for the short run", {})], 0, {})],
        size=25, line_spacing_pts=16)
    rows = [["", "Your costs at that output"],
            ["AVC", "$12"],
            ["ATC", "$15"],
            ["MC", "$10"]]
    _add_styled_table(
        slide, Inches(0.75), Inches(3.45), Inches(4.60), Inches(2.30), rows,
        col_widths=[Inches(1.75), Inches(2.85)],
        row_heights=[Inches(0.62)] + [Inches(0.56)] * 3,
        font_size=20, header_size=19)
    # 2026-08-30 (Nico): the question is what the poll asks, so it goes in a
    # gold box of its own — and the Poll Break badge is drawn AFTER the
    # footer so it sits in front of the rule it straddles.
    _add_takeaway_bar(
        slide, "Should you continue to produce in the short run?",
        top=Inches(6.02), width=Inches(7.60), height=Inches(0.58),
        left=Inches(0.75), fill=GOLD, text_color=NAVY, size=24, bold=True,
        rounded=True, shadow=True)
    _add_media_image(slide, "NV_s36_4_7cb0fc52.jpg",
                     left=Inches(8.85), top=Inches(1.55), width=Inches(4.00))
    _draw_footer(slide, FOOTER_TEXT, page_num)
    _add_pollbreak_badge(slide)
    return slide


def slide_40_coffee_solution(prs, page_num):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_SR)
    _draw_action_title(slide, "Coffee Bean Producer: Solution")
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.15), Inches(1.70), Inches(12.0),
        Inches(1.10),
        [("The output level is already profit-maximizing, so for a price "
          "taker P = MC", 0)],
        size=24, line_spacing_pts=0)
    # 2026-08-30 (Nico): P < AVC is its own dark-red line, introduced by
    # "Thus:", and the four lines sit where he placed them by hand — the
    # two facts close together, then a gap before the conclusion.
    for x, y, h, eq, col, sz in (
            (2.30, 2.29, 0.80,
             _omml_text('MC') + _omml_text(' = 10')
             + _omml_text('     →     ') + _omml_run('P')
             + _omml_text(' = 10'), NAVY, 28),
            (2.30, 3.02, 0.80,
             _omml_text('AVC') + _omml_text(' = 12'), NAVY, 28),
            (2.32, 4.28, 0.57,
             _omml_text('Thus:  ') + _omml_run('P') + _omml_text(' < ')
             + _omml_text('AVC'), DARKRED, 28),
            (2.30, 5.02, 0.80,
             _omml_text('You lose money on every unit — stop production'),
             DARKRED, 30)):
        _add_math_equation(slide, Inches(x), Inches(y), Inches(8.7),
                           Inches(h), eq, size_pt=sz, color=col)
    # hand-set 2026-08-30 (Nico): narrower and taller, pinned at x 2.736,
    # with the clause after the dash on its own line.
    _add_takeaway_bar(
        slide,
        "You will then continue to pay the fixed costs in the short run "
        "\n\u2014 that minimizes the loss",
        top=Inches(5.890), width=Inches(7.380), height=Inches(0.680),
        left=Inches(2.736), fill=GOLD, text_color=NAVY,
        size=19, bold=True, rounded=True, shadow=True)
    _draw_footer(slide, FOOTER_TEXT, page_num)
    _add_pollbreak_badge(slide)
    return slide


def slide_41_sr_summary(prs, page_num):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_SR)
    _draw_action_title(slide,
                       "Summary: Short-Run Decisions for a Price Taker")
    _add_math_equation(
        slide, Inches(2.55), Inches(1.62), Inches(8.2), Inches(0.80),
        _omml_text('Optimal output: ')
        + _omml_sup(_omml_run('Q'), _omml_text('*'))
        + _omml_text(' where  ') + _omml_text('MC') + _omml_text(' = ')
        + _omml_text('MR') + _omml_text(' = ') + _omml_run('P'),
        size_pt=26, color=NAVY)
    # 2026-08-30 (Nico): MW slide 37's format — each rule head sits in a
    # cream box with a 1.75 pt gold border — but WITHOUT the numbered
    # circles; the two rules are a pair, not a sequence.
    for i, (head, subs) in enumerate([
            ("If P ≥ AVC at Q*:  continue to produce",
             ["Profits can still be positive or negative",
              "Positive if P ≥ ATC, negative if P < ATC"]),
            ("If P < AVC at Q*:  stop production",
             ["The loss is then TFC",
              "Producing anyway would make the loss larger than TFC"])]):
        y = Inches(2.55) + i * Inches(2.27)
        head_box = _add_rounded_filled_box(
            slide, Inches(0.70), y, Inches(11.90), Inches(0.72), head,
            fill=CREAM, text_color=NAVY, line=GOLD, size=24,
            corner_pct=0.10)
        head_box.line.width = Pt(1.75)
        _add_hierarchical_bullets(
            slide, Inches(1.30), y + Inches(0.86), Inches(10.9),
            Inches(1.00), [(t, 1) for t in subs],
            size=24, sub_size=22, line_spacing_pts=4)
    _add_ps_pointer(slide, label="Problem Set 3")
    _draw_footer(slide, FOOTER_TEXT, page_num)
    return slide


# ==========================================================================
#  2b · FIRM-LEVEL AND MARKET SUPPLY — slides 42 – 45
# ==========================================================================

def slide_43_supply_curve(prs, page_num):
    """Trace the firm's supply curve out of its MC curve: the two prices
    and the two quantities are the ones the worked example already
    delivered, read straight off the constants."""

    def draw(slide):
        left = SimpleFig(1.30, 6.05, 4.55, 3.70, xmax=820.0, ymax=500.0)
        right = SimpleFig(7.55, 6.05, 4.55, 3.70, xmax=820.0, ymax=500.0)
        for fig, curve_label in ((left, "MC"), (right, "S")):
            _fig_axes(slide, fig, x_title="Q", y_title="Price ($)",
                      label_size=17, titles_at_tip=True)
            _fig_line(slide, fig, (0, mc(0)), (760, mc(760)), color=NAVY,
                      weight_pt=3.0)
            _fig_curve_label(slide, fig, 740, mc(760) + 26, curve_label,
                             size=19)
            for p, q in ((P_LOW, Q_LOW), (P_HIGH, Q_HIGH)):
                _fig_line(slide, fig, (0, p), (q, p), color=GRAY,
                          weight_pt=1.5, dash="dash")
                _fig_line(slide, fig, (q, 0), (q, p), color=GRAY,
                          weight_pt=1.5, dash="dash")
                _fig_ylab(slide, fig, p, _num(p), size=16)
                _xlab_n(slide, fig, q, _num(q), size=16, bold=False, w=0.86)
                _fig_dot(slide, fig, q, p, d=Inches(0.13))
        # AVC only on the left panel: the supply curve is MC above AVC
        _fig_line(slide, left, (0, avc(0)), (760, avc(760)), color=GRAY,
                  weight_pt=2.25)
        _fig_curve_label(slide, left, 740, avc(760) + 26, "AVC", size=17,
                         color=GRAY)
        _add_text(slide, Inches(1.30), Inches(1.58), Inches(4.55),
                  Inches(0.38), "The firm's cost curves", size=20,
                  bold=True, color=NAVY, font="Calibri",
                  align=PP_ALIGN.CENTER)
        _add_text(slide, Inches(7.55), Inches(1.58), Inches(4.55),
                  Inches(0.38), "The firm's supply curve", size=20,
                  bold=True, color=NAVY, font="Calibri",
                  align=PP_ALIGN.CENTER)

    slide = make_diagram_slide(
        prs, page_num, TAG_SUPPLY,
        "The MC Curve Is the Price Taker's Short-Run Supply Curve", draw)
    # 2026-08-30 (Nico): his original wording, with "is" set lower case and
    # underlined rather than shouted in capitals
    bar = _add_mixed_textbox(
        slide, Inches(0.65), Inches(6.42), Inches(8.9), Inches(0.55),
        [("text", "MC curve ", {'size': 20, 'bold': True, 'color': NAVY}),
         ("text", "is", {'size': 20, 'bold': True, 'color': NAVY,
                         'underline': True}),
         ("text", " the firm's short-run supply curve",
          {'size': 20, 'bold': True, 'color': NAVY})],
        align=PP_ALIGN.CENTER)
    bar.fill.solid()
    bar.fill.fore_color.rgb = GOLD
    bar.line.fill.background()
    _apply_picture_style(bar, corner_pct=30)
    bar.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    _yi_badge(slide)
    _set_notes(slide, NOTES[41])
    return slide


# 2026-08-30 (Nico): his original slide 42 shifts MC DOWN in parallel —
# same slope, lower intercept — not a pivot.  High-yield seeds cut the
# cost of every ton by the same amount.
MC1_DROP = 90.0          # dollars per ton, the parallel drop


def mc1(q):
    return mc(q) - MC1_DROP


def slide_44_changing_mc(prs, page_num):
    """High-yield seeds lower marginal cost, so the firm supplies more at
    the same price.  Rebuilt from my original slide 42 (2026-08-30): the
    shift is PARALLEL, the price line is the dark-red d = MR line, and the
    two annotations are his — a down arrow between the curves labelled
    "High-yield seeds", and a q0 -> q1 arrow on the axis."""

    def draw(slide):
        fig = SimpleFig(2.55, 5.90, 7.30, 3.75, xmax=1150.0, ymax=530.0)
        _fig_axes(slide, fig, x_title="q", y_title="P", label_size=18)
        _fig_line(slide, fig, (0, mc(0)), (900, mc(900)), color=NAVY,
                  weight_pt=3.0)
        _fig_curve_label(slide, fig, 880, mc(900) + 30, "MC\u2080", size=19)
        _fig_line(slide, fig, (0, mc1(0)), (1080, mc1(1080)), color=GOLD,
                  weight_pt=3.0)
        _fig_curve_label(slide, fig, 1060, mc1(1080) + 30, "MC\u2081",
                         size=19, color=GOLD)
        _fig_line(slide, fig, (0, P_HIGH), (1120, P_HIGH), color=RED,
                  weight_pt=2.5)
        _add_text(slide, fig.x(1130), fig.y(P_HIGH) - Inches(0.17),
                  Inches(1.60), Inches(0.34), "P = MR", size=18, bold=True,
                  italic=True, color=RED, font="Calibri")
        q0, q1 = Q_HIGH, (P_HIGH + MC1_DROP - B_LIN) / (2 * B_QUAD)
        for q, lab in ((q0, "q\u2080"), (q1, "q\u2081")):
            _fig_line(slide, fig, (q, 0), (q, P_HIGH), color=GRAY,
                      weight_pt=1.5, dash="dash")
            _fig_dot(slide, fig, q, P_HIGH, d=Inches(0.13))
        # his labels sit BESIDE the axis arrow that runs between them
        _add_text(slide, fig.x(q0) - Inches(0.60), fig.y(0) + Inches(0.10),
                  Inches(0.55), Inches(0.34), "q\u2080", size=18, bold=True,
                  italic=True, color=NAVY, font="Calibri",
                  align=PP_ALIGN.RIGHT)
        _add_text(slide, fig.x(q1) + Inches(0.08), fig.y(0) + Inches(0.10),
                  Inches(0.55), Inches(0.34), "q\u2081", size=18, bold=True,
                  italic=True, color=NAVY, font="Calibri")
        _add_arrow(slide, (fig.x(q0) - Inches(0.02), fig.y(0) + Inches(0.27)),
                   (fig.x(q1) + Inches(0.02), fig.y(0) + Inches(0.27)),
                   color=NAVY, weight_pt=2.25, head=True)
        # the down arrow between the two MC curves, at his position
        aq = 620.0
        # dark yellow, matching the MC1 curve it points to (2026-08-30)
        _add_arrow(slide, (fig.x(aq), fig.y(mc(aq) - 18)),
                   (fig.x(aq), fig.y(mc1(aq) + 18)),
                   color=GOLD, weight_pt=3.0, head=True, head_size='lg')
        # hand-set 2026-08-30 (Nico): his wording and position, now a
        # cream bordered card drawn LAST so it sits in FRONT of the MC1
        # curve that runs underneath it.
        _add_convention_box(
            slide, Inches(4.618), Inches(4.125), Inches(1.960),
            Inches(0.640),
            runs=[("High-yield seeds", {'size': 19, 'bold': True,
                                       'color': GOLD}),
                  ("reduce MC", {'size': 19, 'bold': True,
                                 'color': GOLD, 'newline': True})],
            corner_pct=0.14, size=19, border=GOLD,
            align=PP_ALIGN.LEFT, pad_h=0, pad_v=0)

    slide = make_diagram_slide(
        prs, page_num, TAG_SUPPLY,
        "The Effect of Changing Marginal Costs on a Price Taker's Supply",
        draw)
    _set_notes(slide, NOTES[42])
    return slide


def slide_45_market_dynamics(prs, page_num):
    """When EVERY farmer gets the new seeds, market supply shifts, the
    price falls, and the firm ends up between its old output and what it
    would have produced as the only one with the new seeds."""

    def draw(slide):
        mk = SimpleFig(1.05, 6.02, 4.35, 3.35, xmax=11.0, ymax=11.0)
        fm = SimpleFig(7.35, 6.02, 4.45, 3.35, xmax=11.0, ymax=11.0)
        _add_text(slide, Inches(1.05), Inches(1.98), Inches(4.35),
                  Inches(0.38), "Market", size=21, bold=True, color=NAVY,
                  font="Calibri", align=PP_ALIGN.CENTER)
        _add_text(slide, Inches(7.35), Inches(1.98), Inches(4.45),
                  Inches(0.38), "Firm (farmer)", size=21, bold=True,
                  color=NAVY, font="Calibri", align=PP_ALIGN.CENTER)
        for fig in (mk, fm):
            _fig_axes(slide, fig, x_title="Q", y_title="P",
                      label_size=17)

        # market: S0 and S1 = S0 shifted right; D falls
        _fig_line(slide, mk, (1.0, 1.0), (9.4, 9.4), color=NAVY,
                  weight_pt=2.75)
        _fig_curve_label(slide, mk, 9.35, 9.7, "S\u2080", size=18)
        _fig_line(slide, mk, (3.4, 1.0), (10.4, 8.0), color=GOLD,
                  weight_pt=2.75)
        _fig_curve_label(slide, mk, 10.3, 8.4, "S\u2081", size=18, color=GOLD)
        # 2026-08-30 (Nico): the demand curve in dark red
        _fig_line(slide, mk, (1.0, 9.0), (9.6, 0.4), color=RED,
                  weight_pt=2.75)
        _fig_curve_label(slide, mk, 9.5, 0.9, "D", size=18, color=RED)
        # the initial equilibrium (S0 x D) gets a dark-red dot; the new one
        # (S1 x D) stays dark yellow
        for q, p, lab, dotfill in ((5.0, 5.0, "Q\u2080", RED),
                                   (6.2, 3.8, "Q\u2081", GOLD)):
            _fig_line(slide, mk, (0, p), (q, p), color=GRAY,
                      weight_pt=1.5, dash="dash")
            _fig_line(slide, mk, (q, 0), (q, p), color=GRAY,
                      weight_pt=1.5, dash="dash")
            _xlab_n(slide, mk, q, lab, w=0.55)
            _fig_dot(slide, mk, q, p, d=Inches(0.13), fill=dotfill)
        # the supply shift, S0 -> S1 (both lines have slope 1, so the gap
        # is a constant 2.4 in price)
        _add_arrow(slide, (mk.x(8.0), mk.y(8.0 - 0.25)),
                   (mk.x(8.0), mk.y(8.0 - 2.4 + 0.25)),
                   color=GOLD, weight_pt=3.0, head=True, head_size='lg')
        _fig_ylab(slide, mk, 5.0, "P\u2080", size=17, bold=True, color=RED)
        _fig_ylab(slide, mk, 3.8, "P\u2081", size=17, bold=True, color=GOLD)

        # firm: MC0 and MC1, the two price lines, three quantities
        _fig_line(slide, fm, (0, 0), (8.2, 9.84), color=NAVY,
                  weight_pt=2.75)
        _fig_curve_label(slide, fm, 8.05, 10.2, "MC\u2080", size=18)
        # MC1 = MC0 - 2.4, so it meets the x axis at q = 2.0; drawn from
        # there rather than running off below the axis (2026-08-30, Nico)
        # hand-set 2026-08-30 (Nico): MC1 is drawn shorter, ending at
        # q = 9.08 (p = 8.50), with its label pulled in to that end
        _fig_line(slide, fm, (2.0, 0), (9.08, 8.496), color=GOLD,
                  weight_pt=2.75)
        _fig_curve_label(slide, fm, 9.30, 8.51, "MC\u2081", size=18,
                         color=GOLD)
        _fig_line(slide, fm, (0, 5.0), (10.6, 5.0), color=RED,
                  weight_pt=2.25)
        _fig_line(slide, fm, (0, 3.8), (10.6, 3.8), color=GOLD,
                  weight_pt=2.25)
        _add_text(slide, Inches(11.80), fm.y(5.0) - Inches(0.17),
                  Inches(1.50), Inches(0.34), "MR\u2080 = P\u2080", size=15,
                  bold=True, color=RED, font="Calibri")
        _add_text(slide, Inches(11.80), fm.y(3.8) - Inches(0.17),
                  Inches(1.50), Inches(0.34), "MR\u2081 = P\u2081", size=15,
                  bold=True, color=GOLD, font="Calibri")
        # MC1 is MC0 shifted down in parallel, so the quantities move
        # with the intercept, not with a pivot
        for q, p, lab in ((5.0 / 1.2, 5.0, "q₀"),
                          ((3.8 + 2.4) / 1.2, 3.8, "q₂"),
                          ((5.0 + 2.4) / 1.2, 5.0, "q₁")):
            _fig_line(slide, fm, (q, 0), (q, p), color=GRAY,
                      weight_pt=1.25, dash="dash")
            _xlab_n(slide, fm, q, lab)
        # 2026-08-30 (Nico): no dots on the firm panel, and a dark-yellow
        # arrow for the cost shift.  MC0 = 1.2q and MC1 = 1.2q - 2.4, so at
        # q = 7.5 the arrow runs 9.0 -> 6.6, clear of both price lines.
        _add_arrow(slide, (fm.x(7.5), fm.y(9.0 - 0.25)),
                   (fm.x(7.5), fm.y(6.6 + 0.25)),
                   color=GOLD, weight_pt=3.0, head=True, head_size='lg')

    slide = make_diagram_slide(prs, page_num, TAG_SUPPLY,
                               "Firm and Market Dynamics When MC Falls", draw)
    _add_takeaway_bar(
        slide, "Everyone's costs fall, so the price falls too — the farmer "
               "gains less than if the seeds were his alone",
        top=Inches(6.45), width=Inches(11.4), height=Inches(0.55),
        left=(SLIDE_W - Inches(11.4)) // 2, fill=GOLD, text_color=NAVY,
        size=18, bold=True, rounded=True, shadow=True)
    _set_notes(slide, NOTES[43])
    return slide


# ==========================================================================
#  2c · LONG-RUN COMPETITIVE EQUILIBRIUM — slides 46 – 54
# ==========================================================================

def slide_47_long_run(prs, page_num):
    # 2026-08-30 (Nico): the wording and structure of my original slide 45
    # — "What changes?" as the general question, then the TWO things that
    # change as numbered points, each with its own sub-points.
    bullets = [
        ("What changes in the long run?", 0,
         {'bullet_style': 'none', 'bold': True}),
        ("1.   Capital can also be adjusted", 0,
         {'bullet_style': 'none'}),
        ("All costs are considered variable", 1),
        ("Focus on Long Run Average Costs (LAC)", 1),
        ("2.   Market entry and exit", 0, {'bullet_style': 'none'}),
        ("Affect the market price", 1),
    ]
    slide = content_slide(prs, page_num, TAG_LR,
                          "Perfect Competition in the Long Run", bullets,
                          size=28, sub_size=25, line_spacing_pts=18,
                          bullets_width=Inches(7.90),
                          bullets_top=Inches(1.12))
    # 2026-08-30 (Nico): the old / new factory picture on the right — what
    # "capital can also be adjusted" looks like — enlarged and nudged left
    # to the size he set by hand
    _add_media_image(slide, "Old_New_Factory.png",
                     left=Inches(8.37), top=Inches(1.92),
                     width=Inches(4.33))
    _set_notes(slide, NOTES[45])
    return slide


# Long-run curves.  LAC is U-shaped; LMC is derived from it, so LMC cuts
# LAC exactly at LAC's minimum rather than near it:
#     LAC(q) = a(q − q0)² + m      LMC(q) = LAC + q·LAC'(q)
_LR_A, _LR_Q0, _LR_M = 0.25, 5.0, 3.0


def lr_lac(q):
    return _LR_A * (q - _LR_Q0) ** 2 + _LR_M


# 2026-08-30 (Nico): LMC is drawn as a STRAIGHT line, exactly as in my
# original slide 46 — "we want to keep this as simple as possible".  The
# one property that has to hold is that it cuts LAC at LAC's minimum, and
# a line through (Q0, LAC(Q0)) does that for any positive slope.  The
# slope is chosen so the line clears LAC comfortably on the right (they
# would meet again at q = Q0 + slope/_LR_A = 9.6, past the drawn range).
_LR_LMC_SLOPE = 1.15


def lr_lmc(q):
    return _LR_LMC_SLOPE * (q - _LR_Q0) + _LR_M


def lr_q_star(p):
    """Where LMC = p.  Exact, since LMC is now linear."""
    return _LR_Q0 + (p - _LR_M) / _LR_LMC_SLOPE


def slide_48_lr_equilibrium(prs, page_num):
    """Entry pushes the price down from P1, exit pushes it up from P2, and
    both stop at P_LR = min LAC, where LMC = LAC = P."""

    def draw(slide):
        fig = SimpleFig(2.55, 6.10, 7.10, 4.05, xmax=9.6, ymax=9.0)
        _fig_axes(slide, fig, x_title="Q", y_title="P", label_size=18)
        # LAC runs to q = 8.8, where it reaches 6.61 — above the P1 line at
        # 6.0, as Nico asked (2026-08-30)
        LAC_HI = 8.8
        _fig_curve(slide, fig, lr_lac, 1.5, LAC_HI, color=GOLD,
                   weight_pt=3.25, segments=5)
        _fig_curve_label(slide, fig, 8.95, lr_lac(LAC_HI) - 0.30, "LAC",
                         size=19, color=GOLD)
        # LMC: a straight line from just above the axis to past the P1 line
        LMC_LO, LMC_HI = 2.83, 8.30
        _fig_line(slide, fig, (LMC_LO, lr_lmc(LMC_LO)),
                  (LMC_HI, lr_lmc(LMC_HI)), color=NAVY, weight_pt=3.0)
        _fig_curve_label(slide, fig, 8.20, lr_lmc(LMC_HI) + 0.60, "LMC",
                         size=19)

        p_lr = lr_lac(_LR_Q0)
        # the price lines are red (the deck's price colour); P1 and P2 are
        # DASHED, so the long-run equilibrium price stays the solid one
        for p, lab, col, dsh in ((6.0, "P1", RED, "dash"),
                                 (p_lr, "PLR", DARKRED, None),
                                 (1.5, "P2", RED, "dash")):
            q = lr_q_star(p)
            _fig_line(slide, fig, (0, p), (9.2, p), color=col,
                      weight_pt=2.5, dash=dsh)
            _add_text(slide, Inches(9.80), fig.y(p) - Inches(0.17),
                      Inches(2.0), Inches(0.34), "MR = %s" % lab, size=17,
                      bold=True, color=col, font="Calibri")
            _fig_ylab(slide, fig, p, lab, size=17, bold=True)
            _fig_line(slide, fig, (q, 0), (q, p), color=GRAY,
                      weight_pt=1.4, dash="dash")
            _xlab_n(slide, fig, q, lab.replace("P", "Q"))
            _fig_dot(slide, fig, q, p, d=Inches(0.13))

        # 2026-08-30 (Nico): both arrows thicker and dark red, like MW's,
        # with their labels in the same colour so each reads with its arrow
        _add_arrow(slide, (fig.x(7.4), fig.y(5.85)),
                   (fig.x(7.4), fig.y(3.15)), color=DARKRED, weight_pt=4.0,
                   head=True, head_size='lg')
        _add_text(slide, fig.x(7.4) + Inches(0.16),
                  fig.y(4.6) - Inches(0.17),
                  Inches(1.3), Inches(0.34), "Entry", size=17, bold=True,
                  color=DARKRED, font="Calibri")
        _add_arrow(slide, (fig.x(2.6), fig.y(1.65)),
                   (fig.x(2.6), fig.y(2.85)), color=DARKRED, weight_pt=4.0,
                   head=True, head_size='lg')
        _add_text(slide, fig.x(2.6) + Inches(0.16),
                  fig.y(2.3) - Inches(0.17),
                  Inches(1.3), Inches(0.34), "Exit", size=17, bold=True,
                  color=DARKRED, font="Calibri")

    slide = make_diagram_slide(
        prs, page_num, TAG_LR,
        "Long-Run Competitive Equilibrium With Many Identical Producers",
        draw)
    _add_takeaway_bar(
        slide, "At QLR:  P = LAC = LMC, and economic profit is zero",
        top=Inches(6.45), width=Inches(8.2), height=Inches(0.55),
        left=(SLIDE_W - Inches(8.2)) // 2, fill=GOLD, text_color=NAVY,
        size=20, bold=True, rounded=True, shadow=True)
    _set_notes(slide, NOTES[46])
    return slide


def slide_49_chickens(prs, page_num):
    """Cecile Steele: an accidental order of 500 chicks, a large profit, and
    then the entry that competed it away."""
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_LR)
    _draw_action_title(slide,
                       "Entry Follows Profit: The Broiler Chicken Industry")
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.15), Inches(1.70), Inches(6.9),
        Inches(2.90),
        [("In the 1920s the farmer Cecile Steele of Ocean View, Delaware "
          "kept a small flock of chickens for eggs", 0),
         ("In 1923 she ordered 50 chicks and, by accident, received 500. "
          "She kept them and made a sizable profit", 0),
         ("Word of her success spread, and by 1928 hundreds of farmers in "
          "the area were raising chickens for meat", 0)],
        size=23, line_spacing_pts=16)
    _add_media_image(slide, "NV_s47_2_f42d0c16.png",
                     left=Inches(7.55), top=Inches(1.85), width=Inches(4.30))
    _add_text(slide, Inches(7.55), Inches(5.30), Inches(4.30), Inches(0.30),
              "An early Delaware broiler farm. Photo: Vox", size=12,
              italic=True, color=GRAY, font="Calibri",
              align=PP_ALIGN.CENTER)
    _add_takeaway_bar(
        slide, "Profit in a market with low entry barriers → new firms enter",
        top=Inches(6.10), width=Inches(9.0), height=Inches(0.58),
        left=Inches(0.60), fill=GOLD, text_color=NAVY, size=20, bold=True,
        rounded=True, shadow=True)
    _draw_footer(slide, FOOTER_TEXT, page_num)
    _set_notes(slide, NOTES[47])
    return slide


def slide_50_lr_summary(prs, page_num):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_LR)
    _draw_action_title(slide,
                       "Summary: Long-Run Decisions for a Price Taker")
    _add_math_equation(
        slide, Inches(2.55), Inches(1.37), Inches(8.2), Inches(0.80),
        _omml_text('Optimal output: ')
        + _omml_sup(_omml_run('Q'), _omml_text('*'))
        + _omml_text(' where  ') + _omml_text('LMC') + _omml_text(' = ')
        + _omml_text('MR') + _omml_text(' = ') + _omml_run('P'),
        size_pt=26, color=NAVY)
    # 2026-08-30 (Nico): back to the boxed version, at the heights he set
    # by hand — cream box with a 1.75 pt gold border round each rule head,
    # sub-bullets underneath, exactly as slide 45 does it.
    for y_head, y_subs, head, subs in (
            (2.305, 3.165, "If P ≥ LAC at Q*:  stay in the industry",
             ["Positive profits attract entry in the long run"]),
            (4.195, 5.055, "If P < LAC at Q*:  exit, or never enter",
             ["Losses cannot be sustained once capital is adjustable"])):
        head_box = _add_rounded_filled_box(
            slide, Inches(0.70), Inches(y_head), Inches(11.90),
            Inches(0.72), head, fill=CREAM, text_color=NAVY, line=GOLD,
            size=24, corner_pct=0.10)
        head_box.line.width = Pt(1.75)
        _add_hierarchical_bullets(
            slide, Inches(1.30), Inches(y_subs), Inches(10.9),
            Inches(0.55), [(t, 1) for t in subs],
            size=24, sub_size=22, line_spacing_pts=4)
    # the closing definition box is NAVY rather than cream (2026-08-30)
    _add_convention_box(
        slide, Inches(1.75), Inches(5.80), Inches(9.8), Inches(1.10),
        runs=[("Long-run competitive equilibrium:  LMC = P = LAC",
               {'bold': True, 'size': 20, 'color': WHITE}),
              ("\n", {}),
              ("“Normal profits” — economic profit is zero in the long run",
               {'size': 18, 'color': WHITE})],
        fill_rgb=NAVY, border=NAVY,
        corner_pct=0.10, size=20, align=PP_ALIGN.CENTER)
    _draw_footer(slide, FOOTER_TEXT, page_num)
    return slide


def slide_51_drug_market(prs, page_num):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_LR)
    _draw_action_title(slide, "Long-Run Equilibrium: The Drug Market")
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.15), Inches(1.65), Inches(12.4),
        Inches(1.20),
        [("Can we use the toolbox on the market for illegal drugs?", 0),
         ("Assume individual dealers are price takers", 1)],
        size=25, sub_size=23, line_spacing_pts=10)
    _add_text(slide, MARGIN + Inches(0.15), Inches(3.05), Inches(12.4),
              Inches(0.42),
              "The government has two one-time options:", size=25,
              color=NAVY, font="Calibri")
    for i, opt in enumerate(["Arrest the drug dealers",
                             "Arrest or punish the drug users"]):
        y = Inches(3.70) + i * Inches(0.85)
        _add_rounded_filled_box(slide, Inches(1.05), y, Inches(0.55),
                                Inches(0.55), str(i + 1), fill=GOLD,
                                text_color=NAVY, size=24, corner_pct=0.50)
        _add_text(slide, Inches(1.95), y + Inches(0.03), Inches(10.5),
                  Inches(0.50), opt, size=26, color=NAVY, font="Calibri")
    _add_takeaway_bar(
        slide, "Which policy does more to reduce the use of illegal drugs?",
        top=Inches(5.75), width=Inches(9.6), height=Inches(0.58),
        left=(SLIDE_W - Inches(9.6)) // 2, fill=GOLD, text_color=NAVY,
        size=21, bold=True, rounded=True, shadow=True)
    _draw_footer(slide, FOOTER_TEXT, page_num)
    _set_notes(slide, NOTES[49])
    return slide


def slide_53_arrest_dealers(prs, page_num):
    """Deliberately blank below the title — I draw the diagram live in
    class (2026-08-28, Nico).  Chrome and speaker notes only."""
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_LR)
    _draw_action_title(slide, "Arrest Drug Dealers")
    _draw_footer(slide, FOOTER_TEXT, page_num)
    _set_notes(slide, NOTES[51])
    return slide


def slide_54_arrest_users(prs, page_num):
    """Deliberately blank below the title — see slide_53_arrest_dealers."""
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_LR)
    _draw_action_title(slide, "Arrest / Punish Drug Users")
    _draw_footer(slide, FOOTER_TEXT, page_num)
    _set_notes(slide, NOTES[52])
    return slide


# ==========================================================================
#  3 · MARKET DISTORTIONS AND REGULATIONS — slides 55 – 75
# ==========================================================================

def _fig_poly(slide, fig, pts, *, fill=CREAM, line=NAVY, alpha=None):
    """A filled polygon in logical coordinates (surplus triangles, tax
    areas).  Built as a custGeom freeform so it stays editable."""
    dev = [(int(fig.x(x)), int(fig.y(y))) for x, y in pts]
    left = min(p[0] for p in dev)
    top = min(p[1] for p in dev)
    w = max(max(p[0] for p in dev) - left, 1)
    h = max(max(p[1] for p in dev) - top, 1)
    path = ['<a:moveTo><a:pt x="%d" y="%d"/></a:moveTo>'
            % (dev[0][0] - left, dev[0][1] - top)]
    for x, y in dev[1:]:
        path.append('<a:lnTo><a:pt x="%d" y="%d"/></a:lnTo>'
                    % (x - left, y - top))
    path.append('<a:close/>')
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    spPr = shp._element.spPr
    old = spPr.find(qn('a:prstGeom'))
    geom = _H.ET.fromstring(
        '<a:custGeom xmlns:a="%s"><a:avLst/><a:gdLst/><a:ahLst/>'
        '<a:cxnLst/><a:rect l="0" t="0" r="r" b="b"/>'
        '<a:pathLst><a:path w="%d" h="%d">%s</a:path></a:pathLst>'
        '</a:custGeom>' % (_H.A_NS, w, h, "".join(path)))
    old.addnext(geom)
    spPr.remove(old)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1.25)
    shp.shadow.inherit = False
    if alpha is not None:
        sf = spPr.find(qn('a:solidFill'))
        sf.find(qn('a:srgbClr')).append(_H.ET.fromstring(
            '<a:alpha xmlns:a="%s" val="%d"/>' % (_H.A_NS, int(alpha))))
    return shp


def slide_56_distortions_intro(prs, page_num):
    bullets = [
        ("In a competitive market, equilibrium is where supply meets "
         "demand", 0),
        ("Outside forces can push the market away from it", 0),
        ("Price controls — rent control, minimum wages", 1),
        ("Taxes and subsidies", 1),
        ("How do we tell who gains and who loses?", 0),
        ("Consumer and producer surplus", 1),
        ("Deadweight loss", 1),
    ]
    # hand-set 2026-08-30 (Nico): his title, and the bullets narrowed to
    # 6.39" so the competitive-market illustration takes the right half
    slide = content_slide(
        prs, page_num, TAG_DIST,
        "Welfare in a Market Economy and Distortions: Introduction",
        bullets, size=26, sub_size=24, line_spacing_pts=14,
        bullets_width=Inches(6.39))
    _add_text(slide, Inches(7.05), Inches(1.55), Inches(5.70),
              Inches(0.36),
              "A competitive market with many consumers and producers",
              size=15, bold=True, italic=True, color=NAVY, font="Calibri",
              align=PP_ALIGN.CENTER)
    _add_media_image(slide, "Competitve Market Illustration.png",
                     left=Inches(7.05), top=Inches(2.00),
                     width=Inches(5.70))
    return slide


# a shared supply / demand frame for the surplus and price-control slides
def _sd_fig(left=1.90, bottom=6.35, w=4.30, h=4.20):
    return SimpleFig(left, bottom, w, h, xmax=11.0, ymax=11.0)


# The two lines, as endpoints, with the crossing DERIVED from them.  It
# used to be typed in as (3.75, 5.5), which is a point on the SUPPLY line
# and not the intersection at all - so the equilibrium dot, P*, Q* and the
# surplus triangles all sat off the demand curve (2026-08-30, Nico).
SD_D0, SD_D1 = (0.0, 10.0), (9.6, 0.4)      # demand: falls
SD_S0, SD_S1 = (0.0, 1.0), (8.0, 10.6)      # supply: rises


def _sd_cross():
    (dx0, dy0), (dx1, dy1) = SD_D0, SD_D1
    (sx0, sy0), (sx1, sy1) = SD_S0, SD_S1
    md = (dy1 - dy0) / (dx1 - dx0)
    ms = (sy1 - sy0) / (sx1 - sx0)
    q = (sy0 - dy0 - ms * sx0 + md * dx0) / (md - ms)
    return q, dy0 + md * (q - dx0)


SD_QSTAR, SD_PSTAR = _sd_cross()        # = (4.0909, 5.9091)


def _sd_curves(slide, fig, *, d_label="D", s_label="S"):
    """D falls, S rises, and they cross at (SD_QSTAR, SD_PSTAR).

    2026-08-30 (Nico): the demand curve and its label are DARK RED - the
    deck default for a demand curve - and both labels sit where he placed
    them by hand on slide 61.
    """
    _fig_line(slide, fig, SD_D0, SD_D1, color=RED, weight_pt=2.75)
    _fig_curve_label(slide, fig, 9.28, 1.17, d_label, size=19, color=RED)
    _fig_line(slide, fig, SD_S0, SD_S1, color=NAVY, weight_pt=2.75)
    _fig_curve_label(slide, fig, 8.18, 10.75, s_label, size=19)


def _sd_centroid(third):
    """Centre of the surplus triangle whose third vertex is `third` (the
    demand intercept for CS, the supply intercept for PS); the other two
    are (0, P*) and (Q*, P*).  Used to drop the CS / PS label INSIDE it."""
    return ((0.0 + 0.0 + SD_QSTAR) / 3.0,
            (SD_PSTAR + third + SD_PSTAR) / 3.0)


# MW's surplus palette, adopted 2026-08-30 (Nico)
CS_RED = RGBColor(0xC0, 0x20, 0x1B)         # consumer-surplus wash + card
PS_BLUE = RGBColor(0x4E, 0x79, 0xB5)        # producer-surplus wash
PS_BLUE_LINE = RGBColor(0x2E, 0x5A, 0xA8)   # producer-surplus card border


def slide_57_consumer_surplus(prs, page_num):
    def draw(slide):
        fig = _sd_fig(bottom=6.40, h=3.95)
        _fig_axes(slide, fig, x_title="Q", y_title="P", label_size=18)
        # the hypotenuse runs from the demand intercept to (Q*, P*), so it
        # lies exactly ON the demand curve
        _fig_poly(slide, fig,
                  [(0, SD_PSTAR), SD_D0, (SD_QSTAR, SD_PSTAR)],
                  fill=CS_RED, alpha=26000, line=CS_RED)
        _sd_curves(slide, fig)
        _fig_line(slide, fig, (0, SD_PSTAR), (SD_QSTAR, SD_PSTAR),
                  color=GRAY, weight_pt=1.5, dash="dash")
        _fig_line(slide, fig, (SD_QSTAR, 0), (SD_QSTAR, SD_PSTAR),
                  color=GRAY, weight_pt=1.5, dash="dash")
        _fig_ylab(slide, fig, SD_PSTAR, "P*", size=18, bold=True)
        _fig_xlab(slide, fig, SD_QSTAR, "Q*", size=18, bold=True)
        _fig_dot(slide, fig, SD_QSTAR, SD_PSTAR)
        cq, cp = _sd_centroid(SD_D0[1])
        _add_text(slide, fig.x(cq) - Inches(0.45),
                  fig.y(cp) - Inches(0.20),
                  Inches(0.90), Inches(0.40), "CS", size=22, bold=True,
                  color=NAVY, font="Calibri", align=PP_ALIGN.CENTER)
        # hand-set 2026-08-30 (Nico): his position and size, the sentence
        # split onto two paragraphs with 6 pt between them
        _add_convention_box(
            slide, Inches(7.00), Inches(3.34), Inches(5.55), Inches(1.43),
            runs=[("Consumer surplus (CS) is the area below the demand "
                   "curve and above the price line", {'size': 20}),
                  ("What buyers were willing to pay, less what they "
                   "actually paid", {'size': 20, 'newline': True})],
            fill_rgb=CS_RED, border=CS_RED, fill_alpha=26,
            corner_pct=0.10, size=20, pad_h=Inches(0.20),
            pad_v=Inches(0.03), space_before_pts=6)

    slide = make_diagram_slide(prs, page_num, TAG_DIST, "Consumer Surplus",
                               draw)
    _add_text(slide, MARGIN + Inches(0.15), Inches(1.36), Inches(12.4),
              Inches(0.42),
              "CS: the gap between what consumers would pay (the D curve) "
              "and what they do pay (P*)", size=21, color=NAVY,
              font="Calibri")
    _set_notes(slide, NOTES[55])
    return slide


def slide_58_producer_surplus(prs, page_num):
    def draw(slide):
        fig = _sd_fig(bottom=6.40, h=3.95)
        _fig_axes(slide, fig, x_title="Q", y_title="P", label_size=18)
        # the hypotenuse runs from the supply intercept to (Q*, P*), so it
        # lies exactly ON the supply curve
        _fig_poly(slide, fig,
                  [(0, SD_PSTAR), SD_S0, (SD_QSTAR, SD_PSTAR)],
                  fill=PS_BLUE, alpha=34000, line=PS_BLUE)
        _sd_curves(slide, fig)
        _fig_line(slide, fig, (0, SD_PSTAR), (SD_QSTAR, SD_PSTAR),
                  color=GRAY, weight_pt=1.5, dash="dash")
        _fig_line(slide, fig, (SD_QSTAR, 0), (SD_QSTAR, SD_PSTAR),
                  color=GRAY, weight_pt=1.5, dash="dash")
        _fig_ylab(slide, fig, SD_PSTAR, "P*", size=18, bold=True)
        _fig_xlab(slide, fig, SD_QSTAR, "Q*", size=18, bold=True)
        _fig_dot(slide, fig, SD_QSTAR, SD_PSTAR)
        cq, cp = _sd_centroid(SD_S0[1])
        _add_text(slide, fig.x(cq) - Inches(0.45),
                  fig.y(cp) - Inches(0.20),
                  Inches(0.90), Inches(0.40), "PS", size=22, bold=True,
                  color=NAVY, font="Calibri", align=PP_ALIGN.CENTER)
        _add_convention_box(
            slide, Inches(7.25), Inches(3.17), Inches(5.55), Inches(1.52),
            runs=[("Producer surplus (PS) is the area above the supply "
                   "curve and below the price line", {'size': 19}),
                  ("What sellers received, less the least they would have "
                   "accepted", {'size': 19, 'newline': True})],
            fill_rgb=PS_BLUE, border=PS_BLUE_LINE, fill_alpha=34,
            corner_pct=0.10, size=19, pad_h=Inches(0.20),
            pad_v=0, space_before_pts=6)

    slide = make_diagram_slide(prs, page_num, TAG_DIST, "Producer Surplus",
                               draw)
    _add_text(slide, MARGIN + Inches(0.15), Inches(1.36), Inches(12.4),
              Inches(0.42),
              "PS: the gap between the price at which producers would sell "
              "(the S curve) and what they do receive (P*)", size=21,
              color=NAVY, font="Calibri")
    _set_notes(slide, NOTES[56])
    return slide


def slide_59_deadweight_loss(prs, page_num):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_DIST)
    _draw_action_title(slide, "Market Distortions: Deadweight Loss")
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.15), Inches(1.75), Inches(6.47),
        Inches(2.30),
        [("The inefficiency that a market intervention creates", 0),
         ("Mutually beneficial transactions that no longer happen", 0),
         ("Workers willing to work at the market wage and firms willing "
          "to hire them — but a minimum wage stands in the way", 1)],
        size=26, sub_size=23, line_spacing_pts=18)
    # 2026-08-30 (Nico): the welfare definition sits INSIDE the cream box,
    # under the deadweight-loss formula, rather than floating beneath it.
    # 2026-08-30 (Nico): the bullets are narrower, the picture fills the
    # top right, and the two formulas sit in their own cream boxes at the
    # heights he set by hand — the welfare box narrower than the DWL one.
    _add_media_image(slide, "Red Tape.png",
                     left=Inches(7.55), top=Inches(1.45),
                     width=Inches(3.85))
    _add_math_equation(
        slide, Inches(2.04), Inches(5.52), Inches(9.72), Inches(0.84),
        _omml_text('Deadweight loss', color=DARKRED) + _omml_text(' = ')
        + _omml_sub(_omml_text('Welfare'), _omml_text('free market'))
        + _omml_text(' − ')
        + _omml_sub(_omml_text('Welfare'), _omml_text('regulation')),
        size_pt=26, color=NAVY, fill=CREAM, line=NAVY, rounded=True,
        shadow=True)
    _add_math_equation(
        slide, Inches(4.07), Inches(6.46), Inches(5.20), Inches(0.60),
        _omml_text('with   ') + _omml_text('Welfare') + _omml_text(' = ')
        + _omml_text('CS') + _omml_text(' + ') + _omml_text('PS'),
        size_pt=22, color=NAVY, fill=CREAM, line=NAVY, rounded=True,
        shadow=True)
    _draw_footer(slide, FOOTER_TEXT, page_num)
    return slide


def slide_60_sales_tax(prs, page_num):
    """A 10 % sales tax, rebuilt from my original slide 58 and then
    hand-tuned (2026-08-30, Nico): the figure sits lower and further left,
    the y axis is just P, the bracket and the tax label moved, and the
    consequences are a legend of area symbols rather than bullets.

    The tax is on the listed price, so it is AD VALOREM: supply PIVOTS.
    S1(q) = TAX_K * S0(q), with TAX_K chosen so the wedge at the new
    equilibrium is wide enough to carry readable A / B / C / D areas.
    """

    def draw(slide):
        fig = SimpleFig(1.28, 6.44, 5.60, 4.30, xmax=11.0, ymax=11.0)
        _fig_axes(slide, fig, x_title="Q", y_title="P", label_size=18)
        TAX_K = 1.428571

        def s0(q):
            return 2.0 + 0.8 * q

        def s1(q):
            return TAX_K * s0(q)

        def dd(q):
            return 10.0 - q

        q0 = (10.0 - 2.0) / 1.8
        p0 = dd(q0)
        q1 = (10.0 - 2.0 * TAX_K) / (1.0 + 0.8 * TAX_K)
        pb = dd(q1)
        ps = s0(q1)

        # A dark red, C grey (the two transfers); B dark blue, D dark
        # yellow (the two deadweight triangles), so the pair the legend
        # combines is always told apart by colour as well as shape
        _fig_poly(slide, fig, [(0, p0), (0, pb), (q1, pb), (q1, p0)],
                  fill=DARKRED, alpha=28000, line=NAVY)
        _fig_poly(slide, fig, [(0, ps), (0, p0), (q1, p0), (q1, ps)],
                  fill=GRAY, alpha=28000, line=NAVY)
        _fig_poly(slide, fig, [(q1, p0), (q1, pb), (q0, p0)],
                  fill=PS_BLUE_LINE, alpha=28000, line=NAVY)
        _fig_poly(slide, fig, [(q1, ps), (q1, p0), (q0, p0)],
                  fill=DARKYELLOW, alpha=28000, line=NAVY)

        _fig_line(slide, fig, (0, s0(0)), (8.6, s0(8.6)), color=NAVY,
                  weight_pt=2.75)
        _fig_curve_label(slide, fig, 8.70, 9.13, "S", size=19)
        _fig_line(slide, fig, (0, s1(0)), (6.55, s1(6.55)), color=GREEN_DK,
                  weight_pt=2.75)
        _fig_curve_label(slide, fig, 7.07, 11.47, "S\u2019", size=19,
                         color=GREEN_DK)
        _fig_line(slide, fig, (0, dd(0)), (9.6, dd(9.6)), color=RED,
                  weight_pt=2.75)
        _fig_curve_label(slide, fig, 9.5, dd(9.6) + 0.5, "D", size=19,
                         color=RED)

        for p, lab in ((pb, "P_B"), (p0, "P0"), (ps, "P_S")):
            _fig_line(slide, fig, (0, p), (q0 if p == p0 else q1, p),
                      color=GRAY, weight_pt=1.4, dash="dash")
            _fig_ylab(slide, fig, p, lab, size=17, bold=True)
        for q, lab in ((q1, "Q1"), (q0, "Q0")):
            _fig_line(slide, fig, (q, 0), (q, dd(q)), color=GRAY,
                      weight_pt=1.4, dash="dash")
            _xlab_n(slide, fig, q, lab)
        for lx, ly, lab in ((1.5, (p0 + pb) / 2, "A"),
                            (1.5, (ps + p0) / 2, "C"),
                            (q1 + 0.34, p0 + 0.30, "B"),
                            (q1 + 0.34, p0 - 0.42, "D")):
            _add_text(slide, fig.x(lx) - Inches(0.25),
                      fig.y(ly) - Inches(0.17), Inches(0.5), Inches(0.34),
                      lab, size=17, bold=True, color=NAVY, font="Calibri",
                      align=PP_ALIGN.CENTER)

        # hand-set: the bracket at 0.76 and its label at his position
        _fig_vbrace(slide, fig, ps, pb, 0.76, None, color=DARKRED)
        # hand-set 2026-08-30 (Nico): the label is TURNED 90 degrees and
        # set beside the brace, which is how it fits the narrow margin
        lbl = _add_text(slide, Inches(-0.24), Inches(4.06), Inches(1.70),
                        Inches(0.34), "P_B \u2212 P_S = t", size=15,
                        bold=True, color=DARKRED, font="Calibri",
                        align=PP_ALIGN.CENTER)
        lbl.rotation = 270
        qt = 5.2
        _add_arrow(slide, (fig.x(qt), fig.y(s0(qt))),
                   (fig.x(qt), fig.y(s1(qt))), color=GREEN_DK,
                   weight_pt=2.5, head=True, head_both=True)
        _add_text(slide, fig.x(qt) + Inches(0.10),
                  fig.y(8.01) - Inches(0.17),
                  Inches(1.10), Inches(0.34), "tax (t)", size=15,
                  bold=True, color=GREEN_DK, font="Calibri")

        _add_convention_box(
            slide, Inches(7.70), Inches(1.55), Inches(5.15), Inches(1.75),
            runs=[("Price received by the seller (P_S) drops.",
                   {'size': 18, 'bullet': True}),
                  ("Total price paid by the buyer (P_B, inclusive of tax) "
                   "increases, but by less than the tax.",
                   {'size': 18, 'newline': True, 'bullet': True})],
            corner_pct=0.10, size=18, border=DARKRED,
            align=PP_ALIGN.LEFT,
            pad_h=Inches(0.06), pad_v=Inches(0.15))
        # the four consequences as area symbols, laid out as the areas are
        _welfare_rows(
            # 3.42, not his 3.19: with the bullets in, the card above
            # now reaches 3.30 and the first row ran under it
            slide, Inches(8.01), 3.42, Inches(4.95),
            [("h", [(DARKRED, "sq"), (PS_BLUE_LINE, "tri")],
              "The buyers lose A + B"),
             ("h", [(GRAY, "sq"), (DARKYELLOW, "tri_v")],
              "The sellers lose C + D"),
             ("v", [(PS_BLUE_LINE, "tri"), (DARKYELLOW, "tri_v")],
              "Deadweight loss: B + D"),
             ("v", [(DARKRED, "sq"), (GRAY, "sq")],
              "The government gets A + C")],
            pitch=0.50, size=20)

    slide = make_diagram_slide(
        prs, page_num, TAG_DIST,
        "Sales Tax: a 10 % Tax on Top of the Listed Price", draw)
    box = _add_convention_box(
        slide, Inches(7.70), Inches(5.58), Inches(5.15), Inches(0.86),
        runs=[("Taxes on some goods (e.g. cigarettes, gasoline\u2026) can "
               "still make sense due to externalities (more below)",
               {'bold': True, 'size': 18, 'color': WHITE})],
        fill_rgb=PS_BLUE_LINE, border=PS_BLUE_LINE, corner_pct=0.12,
        size=18, align=PP_ALIGN.CENTER, pad_h=Inches(0.20),
        pad_v=Inches(0.03))
    box.line.width = Pt(1.5)
    _add_ps_pointer(slide, label="Problem Set 3", left=Inches(10.60),
                    top=Inches(6.56))
    _set_notes(slide, NOTES[58])
    return slide


def slide_61_tax_incidence(prs, page_num):
    """Adopted from MW slide 65.  In my deck this argument lived only in
    the speaker notes of the sales-tax slide."""
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_DIST)
    _draw_action_title(slide, "Tax Incidence: Who Really Pays?")
    _add_convention_box(
        slide, Inches(1.20), Inches(1.62), Inches(11.0), Inches(1.05),
        prefix="Tax incidence:  ",
        body="which side of the market — buyers or sellers — ends up "
             "bearing the burden of the tax",
        corner_pct=0.10, size=21)
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.35), Inches(3.05), Inches(12.2),
        Inches(2.85),
        [("Who hands the money to the government does not matter", 0),
         ("The incidence is the same whether the tax is collected from "
          "the buyer or from the seller", 1),
         ("What matters is the elasticity of supply and demand", 0),
         ("The side that responds less to price bears more of the tax", 1)],
        size=26, sub_size=23, line_spacing_pts=16)
    _add_takeaway_bar(
        slide, "The less price-responsive side of the market pays the "
               "larger share",
        top=Inches(6.20), width=Inches(9.8), height=Inches(0.58),
        left=(SLIDE_W - Inches(9.8)) // 2, fill=GOLD, text_color=NAVY,
        size=20, bold=True, rounded=True, shadow=True)
    _draw_footer(slide, FOOTER_TEXT, page_num)
    return slide


def slide_62_minwage_map(prs, page_num):
    """Adopted from MW slide 53; the map is re-fetched from Wikimedia
    Commons rather than lifted, and is the July 2026 edition."""
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_DIST)
    _draw_action_title(slide, "The Minimum Wage: A Popular Policy Tool")
    # 2026-08-30 (Nico): MW's own figure, so the colour scheme is exactly
    # theirs rather than the Commons default
    _add_media_image(slide, "MW_minwage_map.jpg",
                     left=Inches(2.70), top=Inches(1.50),
                     width=Inches(7.93), height=Inches(5.29))
    _add_text(slide, Inches(2.70), Inches(6.86), Inches(7.93), Inches(0.28),
              "Source: Wikipedia / US Department of Labor",
              size=11, italic=True, color=GRAY, font="Calibri",
              align=PP_ALIGN.CENTER)
    _draw_footer(slide, FOOTER_TEXT, page_num)
    return slide


def slide_63_min_wage(prs, page_num):
    """The labour-market price floor, with the circumvention card adopted
    from MW slide 52."""

    def draw(slide):
        fig = SimpleFig(1.30, 6.15, 4.55, 4.15, xmax=11.0, ymax=11.0)
        _fig_axes(slide, fig, x_title="Labor", y_title="Wage",
                  label_size=18)
        w_star, l_star = 5.5, 3.75
        w_min = 7.0
        l_d, l_s = (10.0 - w_min) / 1.2, (w_min - 1.0) / 1.2

        # 2026-08-30 (Nico): A dark red, B grey, C dark green, all washed
        # back so the curves stay readable through them
        _fig_poly(slide, fig, [(0, w_star), (0, w_min), (l_d, w_min),
                               (l_d, w_star)], fill=DARKRED, alpha=28000,
                  line=NAVY)
        _fig_poly(slide, fig, [(l_d, w_star), (l_d, w_min),
                               (l_star, w_star)], fill=GRAY, alpha=28000,
                  line=NAVY)
        _fig_poly(slide, fig, [(l_d, w_star), (l_star, w_star),
                               (l_d, (1.0 + 1.2 * l_d))], fill=GREEN_DK,
                  alpha=28000, line=NAVY)

        _fig_line(slide, fig, (0.0, 10.0), (8.0, 0.4), color=RED,
                  weight_pt=2.75)
        _fig_curve_label(slide, fig, 8.05, 1.30, "D_Labor", size=17,
                         color=RED)
        _fig_line(slide, fig, (0.0, 1.0), (8.0, 10.6), color=NAVY,
                  weight_pt=2.75)
        _fig_curve_label(slide, fig, 8.05, 11.55, "S_Labor", size=17)
        _fig_line(slide, fig, (0, w_star), (l_star, w_star), color=GRAY,
                  weight_pt=1.4, dash="dash")
        _fig_line(slide, fig, (l_star, 0), (l_star, w_star), color=GRAY,
                  weight_pt=1.4, dash="dash")
        _fig_ylab(slide, fig, w_star, "w*", size=17, bold=True)
        _xlab_n(slide, fig, l_star, "L*")
        # the wage floor runs past L_S, so it reads as a policy line
        _fig_line(slide, fig, (0, w_min), (l_s + 1.7, w_min),
                  color=DARKRED, weight_pt=2.75)
        _fig_ylab(slide, fig, w_min, "wmin", size=17, bold=True)
        for l, lab in ((l_d, "L_D"), (l_s, "L_S")):
            _fig_line(slide, fig, (l, 0), (l, w_min), color=GRAY,
                      weight_pt=1.4, dash="dash")
            _xlab_n(slide, fig, l, lab)
        # 2026-08-30 (Nico): the unemployment gap is an under-brace below
        # the axis, as in my original slide, not an arrow across the plot
        _fig_underbrace(slide, fig, l_d, l_s, 6.61,
                        "L_S − L_D = unemployment", color=DARKRED,
                        size=14)
        for lx, ly, lab in ((1.2, 6.25, "A"), (l_d + 0.30, 6.15, "B"),
                            (l_d + 0.30, 4.95, "C")):
            _add_text(slide, fig.x(lx) - Inches(0.25),
                      fig.y(ly) - Inches(0.17), Inches(0.5), Inches(0.34),
                      lab, size=17, bold=True, color=NAVY, font="Calibri",
                      align=PP_ALIGN.CENTER)

        _add_text(slide, Inches(6.55), Inches(1.62), Inches(6.35),
                  Inches(0.40), "Welfare effects", size=21, bold=True,
                  color=NAVY, font="Calibri")
        _welfare_rows(
            slide, Inches(6.70), 2.20, Inches(6.20),
            # the marks are laid out as the areas are: A and B sit side by
            # side in the graph, B and C sit one above the other
            [("h", [(DARKRED, "sq")],
              "Some workers win A — they are paid more"),
             ("h", [(GREEN_DK, "tri_v")],
              "Some workers lose C — they are not hired"),
             ("h", [(DARKRED, "sq"), (GRAY, "tri")], "Firms lose A + B"),
             ("v", [(GRAY, "tri"), (GREEN_DK, "tri_v")],
              "Deadweight loss: B + C")])
        # 2026-08-30 (Nico): the heading carries NO bullet; the three ways
        # round the wage floor each carry one.
        cbox = _add_rounded_filled_box(
            slide, Inches(8.10), Inches(5.04), Inches(4.45), Inches(1.73),
            "", fill=CREAM, text_color=NAVY, line=NAVY, size=19,
            corner_pct=0.10)
        _add_drop_shadow(cbox)
        _add_hierarchical_bullets(
            slide, Inches(8.28), Inches(5.18), Inches(4.13), Inches(1.48),
            [("Circumvention of minimum wages?", 0,
              {'bullet_style': 'none', 'bold': True, 'size': 20}),
             ("Underground labor market", 1),
             ("Cutting non-wage benefits", 1),
             ("Unpaid internships", 1)],
            size=18, sub_size=18, line_spacing_pts=6)

    slide = make_diagram_slide(prs, page_num, TAG_DIST,
                               "Price Floor: The Minimum Wage", draw)
    return slide


def slide_64_minwage_evidence(prs, page_num):
    """My evidence bullets, with the two headline clippings and the
    three-perspective discussion prompt adopted from MW slide 54."""
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_DIST)
    # 2026-08-30 (Nico): he prefixed the title with "Podcast:"
    _draw_action_title(
        slide, "Podcast: What Does the Evidence Say on the Minimum Wage?")
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.15), Inches(1.55), Inches(6.30),
        Inches(3.95),
        # 2026-08-30 (Nico): the wording and levels of my original slide 59
        [("What is the expected effect of raising the minimum wage on "
          "employment?", 0),
         ("According to economic theory: Negative", 1),
         ("Evidence: Mixed results", 0),
         ("Employment does not seem to decline, BUT:", 1),
         ("Hours worked decline (Seattle study)", 1),
         ("Experienced workers: moderate decline in hours, so total "
          "income ↑", 2),
         ("Low-skilled, inexperienced workers: substantial decline in "
          "hours and thus income ↓", 2),
         ("New workers: harder to find first job", 2),
         ("What are possible long-run effects of a (high) minimum wage?",
          0)],
        size=20, sub_size=18, line_spacing_pts=7)
    # 2026-08-30 (Nico): the "How much is too much" card is deleted, which
    # frees the bottom-left corner
    # hand-set 2026-08-30 (Nico): both clippings moved up, and the photo
    # enlarged with its note directly underneath (the grouping pass then
    # binds the two into one object)
    _add_media_image(slide, "mw_headline_krueger_nyt_2015.png",
                     left=Inches(6.92), top=Inches(1.44), width=Inches(5.90),
                     rounded=False, shadow=True)
    _add_media_image(slide, "mw_headline_vox_2017.png",
                     left=Inches(6.95), top=Inches(2.97), width=Inches(5.90),
                     rounded=False, shadow=True)
    _add_media_image(slide, "NV_s59_3_3e55694d.png",
                     left=Inches(6.95), top=Inches(4.34),
                     width=Inches(3.57), height=Inches(2.42))
    _add_text(slide, Inches(7.05), Inches(6.79), Inches(3.39), Inches(0.30),
              "One way firms respond to a higher wage floor",
              size=12, italic=True, color=GRAY, font="Calibri",
              align=PP_ALIGN.CENTER)
    _draw_footer(slide, FOOTER_TEXT, page_num)
    _set_notes(slide, NOTES[59])
    return slide


def slide_65_guess_rent(prs, page_num):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_DIST)
    _draw_action_title(slide, "Guess the Rent!")
    # hand-set 2026-08-30 (Nico): the photo enlarged to 6.38 x 5.41 and the
    # listing details at 28 / 24 pt, with the last two lines as sub-bullets
    # (the structure of my original slide 61)
    _add_media_image(slide, "NV_s61_6_fb7ae0ee.png",
                     left=Inches(0.28), top=Inches(1.54),
                     width=Inches(6.38), height=Inches(5.41))
    _add_hierarchical_bullets(
        slide, Inches(7.11), Inches(1.39), Inches(5.83), Inches(2.80),
        [("1 bed, 1 bath", 0),
         ("900 sqft", 0),
         ("401 San Vicente Blvd., Santa Monica", 0),
         ("5-minute walk to the beach", 1),
         ("Best public school district in West LA", 1)],
        size=28, sub_size=24, line_spacing_pts=14)
    # 2026-08-30 (Nico): the location map, with the map's own black pin
    # covered by a larger dark-red dot (the pin sits at 51.6 % / 28.3 % of
    # the image, measured off the file)
    MAP_X, MAP_Y, MAP_W, MAP_H = 6.89, 4.13, 3.63, 2.92
    _add_media_image(slide, "NV_flat_map.png",
                     left=Inches(MAP_X), top=Inches(MAP_Y),
                     width=Inches(MAP_W), height=Inches(MAP_H))
    dot_d = 0.24
    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        int(Inches(MAP_X + 0.5157 * MAP_W - dot_d / 2)),
        int(Inches(MAP_Y + 0.2831 * MAP_H - dot_d / 2)),
        int(Inches(dot_d)), int(Inches(dot_d)))
    dot.fill.solid()
    dot.fill.fore_color.rgb = DARKRED
    dot.line.color.rgb = DARKRED
    dot.shadow.inherit = False
    _draw_footer(slide, FOOTER_TEXT, page_num)
    _add_pollbreak_badge(slide)
    return slide


def slide_67_the_rent_is(prs, page_num):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_DIST)
    _draw_action_title(slide, "The Rent Is…")
    # 2026-08-30 (Nico): he replaced my three-part composition with MW's
    # single screenshot and drew a red box round the figure that matters.
    _add_media_image(slide, "MW_rent_is.png",
                     left=Inches(1.67), top=Inches(2.25),
                     width=Inches(10.00), height=Inches(2.95),
                     rounded=False, shadow=False)
    box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(3.43), Inches(3.93),
        Inches(6.57), Inches(0.46))
    box.fill.background()
    box.line.color.rgb = RGBColor(0xFF, 0x00, 0x00)
    box.line.width = Pt(2.0)
    box.shadow.inherit = False
    _draw_footer(slide, FOOTER_TEXT, page_num)
    _add_pollbreak_badge(slide)
    _set_notes(slide, NOTES[63])
    return slide


def slide_68_rent_control(prs, page_num):
    def draw(slide):
        fig = SimpleFig(1.30, 6.15, 4.55, 4.15, xmax=11.0, ymax=11.0)
        _fig_axes(slide, fig, x_title="Q", y_title="P", label_size=18)
        p_max = 3.5
        q_s, q_d = (p_max - 1.0) / 1.2, (10.0 - p_max) / 1.2

        # 2026-08-30 (Nico): A dark red, B grey, C dark green, washed back
        _fig_poly(slide, fig, [(0, p_max), (0, SD_PSTAR),
                               (q_s, SD_PSTAR), (q_s, p_max)],
                  fill=DARKRED, alpha=28000, line=NAVY)
        _fig_poly(slide, fig, [(q_s, SD_PSTAR), (SD_QSTAR, SD_PSTAR),
                               (q_s, 10.0 - q_s)], fill=GRAY, alpha=28000,
                  line=NAVY)
        _fig_poly(slide, fig, [(q_s, SD_PSTAR), (SD_QSTAR, SD_PSTAR),
                               (q_s, p_max)], fill=GREEN_DK, alpha=28000,
                  line=NAVY)

        _sd_curves(slide, fig)
        _fig_line(slide, fig, (0, SD_PSTAR), (SD_QSTAR, SD_PSTAR),
                  color=GRAY, weight_pt=1.4, dash="dash")
        _fig_line(slide, fig, (SD_QSTAR, 0), (SD_QSTAR, SD_PSTAR),
                  color=GRAY, weight_pt=1.4, dash="dash")
        _fig_ylab(slide, fig, SD_PSTAR, "P*", size=17, bold=True)
        _xlab_n(slide, fig, SD_QSTAR, "Q*")
        # the ceiling runs past Q_D so the line reads as a policy, not as a
        # segment (2026-08-30, Nico); Q_D is where it meets demand
        _fig_line(slide, fig, (0, p_max), (q_d + 3.1, p_max),
                  color=DARKRED, weight_pt=2.75)
        _fig_ylab(slide, fig, p_max, "Pmax", size=17, bold=True)
        for q, lab in ((q_s, "Q_S"), (q_d, "Q_D")):
            _fig_line(slide, fig, (q, 0), (q, p_max), color=GRAY,
                      weight_pt=1.4, dash="dash")
            _xlab_n(slide, fig, q, lab)
        # 2026-08-30 (Nico): MW marks the shortage with an under-brace
        # below the axis rather than an arrow drawn across the figure
        _fig_underbrace(slide, fig, q_s, q_d, 6.61,
                        "Q_D − Q_S = shortage", color=DARKRED, size=14)
        for lx, ly, lab in ((1.0, 4.45, "A"), (q_s + 0.32, 6.15, "B"),
                            (q_s + 0.32, 4.95, "C")):
            _add_text(slide, fig.x(lx) - Inches(0.25),
                      fig.y(ly) - Inches(0.17), Inches(0.5), Inches(0.34),
                      lab, size=17, bold=True, color=NAVY, font="Calibri",
                      align=PP_ALIGN.CENTER)

        _add_text(slide, Inches(6.55), Inches(1.75), Inches(6.35),
                  Inches(0.40), "Welfare effects", size=21, bold=True,
                  color=NAVY, font="Calibri")
        _welfare_rows(
            slide, Inches(6.70), 2.35, Inches(6.20),
            # shortened 2026-08-30 so the row fits on one line (5.88" ->
            # 4.72" at 18 pt, against a 5.58" column)
            # A sits above C in the graph and B above C too, so both
            # combinations stack rather than sitting side by side
            [("h", [(DARKRED, "sq")],
              "Some renters win A — those who keep a place"),
             ("h", [(GRAY, "tri")],
              "Some renters lose B — they cannot find one"),
             # 2026-08-30 (Nico): A and C sit side by side in the graph
             ("h", [(DARKRED, "sq"), (GREEN_DK, "tri_v")],
              "Landlords lose A + C"),
             ("v", [(GRAY, "tri"), (GREEN_DK, "tri_v")],
              "Deadweight loss: B + C")],
            pitch=0.55, size=18, text_dx=0.62)

    slide = make_diagram_slide(prs, page_num, TAG_DIST,
                               "Price Ceiling: Rent Control", draw)
    return slide


def slide_69_prop33(prs, page_num):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_DIST)
    _draw_action_title(slide, "Vote on Prop 33 in California, 2024")
    _add_media_image(slide, "NV_s65_5_2f36cabb.png",
                     left=Inches(3.55), top=Inches(1.55),
                     height=Inches(2.75), width=Inches(6.20))
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.35), Inches(4.55), Inches(11.4),
        Inches(1.90),
        [("California limits how far cities may go with rent control; "
          "Prop 33 would have removed those limits", 0),
         ("Cities could then have controlled rents on any housing — "
          "single-family homes, new apartments, and new tenants", 0),
         ("It did not pass", 0, {'bold': True})],
        size=21, line_spacing_pts=10)
    _add_text(slide, Inches(0.45), Inches(6.62), Inches(8.4), Inches(0.28),
              "Source: calmatters.org, California voter guide 2024",
              size=11, italic=True, color=GRAY, font="Calibri")
    _draw_footer(slide, FOOTER_TEXT, page_num)
    _add_pollbreak_badge(slide)
    _set_notes(slide, NOTES[65])
    return slide


def slide_71_landlord_reaction(prs, page_num):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_DIST)
    _draw_action_title(slide, "How Do Landlords React to Rent Control?")
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.15), Inches(1.53), Inches(4.36),
        Inches(2.11),
        # hand-set 2026-08-30 (Nico): his wording, tighter spacing, both
        # pictures enlarged and repositioned, and the captions cut
        [("Non-price adjustments (e.g. neglect maintenance)", 0),
         ("Buy the tenants out (legal)", 0),
         ("Harass the tenants (illegal)", 0)],
        size=25, line_spacing_pts=12)
    _add_media_image(slide, "NV_s67_5_9b9eb200.jpg",
                     left=Inches(2.08), top=Inches(3.71),
                     width=Inches(4.42), height=Inches(3.31))
    _add_media_image(slide, "NV_s67_4_868e46cb.jpg",
                     left=Inches(8.40), top=Inches(1.41),
                     width=Inches(4.24), height=Inches(5.66))
    _draw_footer(slide, FOOTER_TEXT, page_num)
    return slide


def slide_72_cambridge_map(prs, page_num):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_DIST)
    _draw_action_title(slide, "Evidence on Rent Control: Cambridge, MA")
    _add_text(slide, MARGIN + Inches(0.15), Inches(1.52), Inches(12.4),
              Inches(0.42),
              "A natural experiment: rent control in Cambridge ended "
              "suddenly in 1995", size=23, color=NAVY, font="Calibri")
    # hand-set 2026-08-30 (Nico): the map enlarged and its source line
    # moved down with it
    _add_media_image(slide, "NV_s68_5_5777e352.png",
                     left=Inches(2.78), top=Inches(2.05),
                     width=Inches(7.68), height=Inches(4.71))
    _add_text(slide, Inches(3.35), Inches(6.86), Inches(6.65), Inches(0.30),
              "Source: Autor, Palmer and Pathak (2014)", size=12,
              italic=True, color=GRAY, font="Calibri",
              align=PP_ALIGN.CENTER)
    _draw_footer(slide, FOOTER_TEXT, page_num)
    # the paper my original slide linked to
    _set_notes(slide, (NOTES[68] + "\n\nSource: Autor, Palmer and Pathak "
                       "(2014), \u201cHousing Market Spillovers: Evidence "
                       "from the End of Rent Control in Cambridge, "
                       "Massachusetts\u201d \u2014 "
                       "https://economics.mit.edu/files/9774").strip())
    return slide


def slide_73_rent_questions(prs, page_num):
    """2026-08-30 (Nico): the wording of my original slide 69, verbatim."""
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_DIST)
    _draw_action_title(slide, "What Happened When Rent Control Ended?")
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.35), Inches(1.90), Inches(12.0),
        Inches(3.60),
        [("\u201cNatural\u201d experiment: sudden end of rent control in "
          "Cambridge, MA in 1995", 0),
         ("What do you think happened to:", 0),
         ("Rental price for previously rent-controlled units?", 1),
         ("Investment in previously rent-controlled units?", 1),
         ("Price/investment for nearby never-controlled units?", 1)],
        size=27, sub_size=24, line_spacing_pts=18)
    _draw_footer(slide, FOOTER_TEXT, page_num)
    _add_pollbreak_badge(slide)
    return slide


def slide_74_rent_results(prs, page_num):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_DIST)
    _draw_action_title(slide, "Evidence on Rent Control: Results")
    # 2026-08-30 (Nico): the headers are the wording of my original slide
    # 70, and each panel carries the red trend line I drew there — mapped
    # onto this layout by the line's position RELATIVE to its own picture,
    # so it sits on the same part of each chart.  All three rise.
    panels = [
        # hand-set 2026-08-30 (Nico): his panel sizes, header boxes and
        # trend-line positions, the lines again stored as fractions of
        # their own picture so they travel with it
        ("NV_s70_2_95e7216e.png", 0.22, 1.96, 4.52, 2.82,
         "Impact on the price of decontrolled units",
         (0.08, 1.65, 4.80, 0.56),
         (0.4173, 0.2273, 0.9151, 0.4472)),
        ("NV_s70_4_5ecaba44.png", 5.29, 1.95, 3.75, 2.76,
         "Impact on the price of never-controlled units",
         (5.00, 1.60, 4.20, 0.56),
         (0.4307, 0.4301, 0.9640, 0.6290)),
        ("NV_s70_6_08feed5a.png", 9.59, 1.97, 3.38, 3.92,
         "Impact on investment activities at decontrolled units ($1000s)",
         (9.65, 1.47, 3.45, 0.47),
         (0.3130, 0.3395, 0.6843, 0.5140)),
    ]
    for fn, x, y, w, h, cap, hd, ln in panels:
        _add_text(slide, Inches(hd[0]), Inches(hd[1]), Inches(hd[2]),
                  Inches(hd[3]), cap, size=14, bold=True,
                  italic=True, color=NAVY, font="Calibri",
                  align=PP_ALIGN.CENTER)
        _add_media_image(slide, fn, left=Inches(x), top=Inches(y),
                         width=Inches(w), height=Inches(h),
                         rounded=False, shadow=True)
        fx0, fy0, fx1, fy1 = ln
        _add_arrow(slide,
                   (Inches(x + fx0 * w), Inches(y + fy1 * h)),
                   (Inches(x + fx1 * w), Inches(y + fy0 * h)),
                   color=RED, weight_pt=2.25, head=False)
    _add_takeaway_bar(
        slide, "Rents rose, and so did investment — including at the units "
               "next door",
        top=Inches(6.33), width=Inches(9.8), height=Inches(0.55),
        left=(SLIDE_W - Inches(9.8)) // 2, fill=PS_BLUE_LINE,
        text_color=WHITE,
        size=19, bold=True, rounded=True, shadow=True)
    _draw_footer(slide, FOOTER_TEXT, page_num)
    # the source moves off the slide and into the notes (2026-08-30, Nico)
    _set_notes(slide, (NOTES.get(70, "") + "\n\nSource of all three "
                       "panels: Autor, Palmer and Pathak (2014), "
                       "\u201cHousing Market Spillovers: Evidence from the "
                       "End of Rent Control in Cambridge, "
                       "Massachusetts\u201d.").strip())
    return slide


def slide_75_argentina(prs, page_num):
    """Adopted from MW slide 63: a second, very recent natural experiment
    beside Cambridge 1995.  The magnitudes are press-reported, and the
    slide says so."""
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_DIST)
    _draw_action_title(slide, "Argentina Scrapped Its Rent Controls")
    _add_media_image(slide, "mw_headline_wsj_argentina.png",
                     left=Inches(1.55), top=Inches(1.65), width=Inches(10.20),
                     rounded=False, shadow=True)
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.55), Inches(3.90), Inches(11.6),
        Inches(1.90),
        [("Argentina repealed its rent-control law at the end of 2023", 0),
         ("Reported since: many more flats listed, and rents lower "
          "in real terms", 0),
         ("Not everyone gains — sitting tenants lost their protection", 0)],
        size=24, line_spacing_pts=16)
    # 2026-08-30 (Nico): the caveat moves off the slide into the notes
    _draw_footer(slide, FOOTER_TEXT, page_num)
    _set_notes(slide, "Magnitudes are as reported in the press, not from "
                      "a study. Source: The Wall Street Journal, "
                      "“Argentina Scrapped Its Rent Controls. Now the "
                      "Market Is Thriving.”")
    return slide


# ==========================================================================
#  4 · EXTERNALITIES — slides 76 – 83
# ==========================================================================

def slide_77_externalities_def(prs, page_num):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_EXT)
    _draw_action_title(slide, "Externalities: Definitions")
    _add_convention_box(
        slide, Inches(1.00), Inches(1.60), Inches(11.3), Inches(1.00),
        prefix="Externality:  ",
        body="A cost or a benefit that falls on someone who is not part "
             "of the transaction",
        corner_pct=0.10, size=22)
    _add_hierarchical_bullets(
        slide, MARGIN + Inches(0.35), Inches(3.00), Inches(12.2),
        Inches(3.30),
        # 2026-08-30 (Nico): the term in bold concept blue, then "A", then
        # the head noun bold AND underlined, then the rest plain
        [([("Negative externality", {'bold': True, 'color': BLUE_PED}),
           (" — A ", {}),
           ("cost", {'bold': True, 'underline': True}),
           (" imposed on an outsider", {})], 0, {}),
         ("Air pollution from coal-fired power plants", 1),
         ("The external cost of gasoline is about $1.50 per gallon", 1),
         ([("Positive externality", {'bold': True, 'color': BLUE_PED}),
           (" — A ", {}),
           ("benefit", {'bold': True, 'underline': True}),
           (" conferred on an outsider", {})], 0, {}),
         ("Your neighbor's well-kept front yard is a pleasure for you "
          "to look at too", 1)],
        size=24, sub_size=22, line_spacing_pts=16)
    _draw_footer(slide, FOOTER_TEXT, page_num)
    return slide


def slide_78_correcting(prs, page_num):
    """Adopted from MW slide 67, with one worked example for each kind of
    quota (2026-08-28, Nico)."""
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_EXT)
    _draw_action_title(slide, "Correcting for Externalities")
    _add_text(slide, MARGIN + Inches(0.15), Inches(1.50), Inches(8.30),
              Inches(0.75),
              "The aim of the intervention: move the market toward the "
              "quantity that would be chosen if the outside cost or "
              "benefit were counted", size=21, color=NAVY, font="Calibri")
    _add_rounded_filled_box(
        slide, Inches(0.55), Inches(2.55), Inches(8.30), Inches(0.55),
        "Price mechanisms:  taxes and subsidies", fill=NAVY,
        text_color=WHITE, size=20, corner_pct=0.12)
    _add_hierarchical_bullets(
        slide, Inches(0.85), Inches(3.20), Inches(8.00), Inches(0.90),
        [("Tax the activity that imposes the cost, subsidize the one "
          "that confers the benefit", 0)],
        size=20, line_spacing_pts=0)
    _add_rounded_filled_box(
        slide, Inches(0.55), Inches(4.15), Inches(8.30), Inches(0.55),
        "Quantity mechanisms:  quotas and mandates", fill=NAVY,
        text_color=WHITE, size=20, corner_pct=0.12)
    _add_hierarchical_bullets(
        slide, Inches(0.85), Inches(4.85), Inches(8.00), Inches(1.60),
        [([("Negative externality — cap it:", {'bold': True}),
           (" a limit on the tonnes of sulphur dioxide a power plant may "
            "emit in a year", {})], 0, {}),
         ([("Positive externality — require it:", {'bold': True}),
           (" childhood vaccination as a condition of school entry", {})],
          0, {})],
        size=19, line_spacing_pts=12)
    _add_media_image(slide, "NV_s73_4_e69d8d25.png",
                     left=Inches(9.35), top=Inches(1.55), width=Inches(3.20))
    _draw_footer(slide, FOOTER_TEXT, page_num)
    return slide


# The gasoline externality figure, shared by slides 79 and 80.
EMC = 1.5


def _gas_d(q):
    return 8.0 - 0.6 * q


def _gas_mci(q):
    return 1.0 + 0.55 * q


def _gas_smc(q):
    return _gas_mci(q) + EMC


GAS_Q_MARKET = 7.0 / 1.15          # 8 − 0.6q = 1 + 0.55q
GAS_Q_EXT = 5.5 / 1.15             # 8 − 0.6q = 2.5 + 0.55q


def _draw_gas_fig(slide, *, with_tax):
    fig = SimpleFig(2.30, 6.35, 6.15, 4.35, xmax=9.4, ymax=9.0)
    _fig_axes(slide, fig, x_title="Gallons of gasoline per year",
              y_title="$ / gallon", label_size=17)
    _fig_line(slide, fig, (0, _gas_d(0)), (8.8, _gas_d(8.8)), color=RED,
              weight_pt=2.75)
    _add_text(slide, Inches(7.71 if with_tax else 8.07),
              Inches(4.54 if with_tax else 4.64), Inches(0.98),
              Inches(0.32), "Demand", size=18, bold=True, italic=True,
              color=RED, font="Calibri")
    _fig_line(slide, fig, (0, _gas_mci(0)), (8.8, _gas_mci(8.8)),
              color=NAVY, weight_pt=2.75)
    _fig_line(slide, fig, (0, _gas_smc(0)), (8.0, _gas_smc(8.0)),
              color=GOLD, weight_pt=2.75)
    # 2026-08-30 (Nico): MW's labelling — a price label on the y axis for
    # each equilibrium as well as the quantity on the x axis, everything
    # with real subscripts.
    for q, p, lab, col in (
            (GAS_Q_MARKET, _gas_d(GAS_Q_MARKET), "Market", NAVY),
            (GAS_Q_EXT, _gas_d(GAS_Q_EXT), "Ext", GOLD)):
        _fig_line(slide, fig, (0, p), (q, p), color=GRAY, weight_pt=1.4,
                  dash="dash")
        _fig_line(slide, fig, (q, 0), (q, p), color=GRAY, weight_pt=1.4,
                  dash="dash")
        _xlab_n(slide, fig, q, "Q_" + lab, size=16)
        # on the tax slide the efficient price is labelled P_E^C instead,
        # so its plain P_Ext tick is suppressed here (2026-08-30, Nico)
        if not (with_tax and lab == "Ext"):
            _fig_ylab(slide, fig, p, "P_" + lab, size=16, bold=True)
        _fig_dot(slide, fig, q, p, d=Inches(0.13), fill=col)
    # the external marginal cost itself: a flat $1.50 per gallon, drawn as
    # a thick black line with a soft shade, as MW draws it
    emc = _fig_line(slide, fig, (0, 1.5), (8.8, 1.5),
                    color=EMC_PURPLE, weight_pt=3.5)
    _add_drop_shadow(emc)
    _fig_ylab(slide, fig, 1.5, "1.5", size=16, bold=True)
    # hand-set 2026-08-30 (Nico): on the wedge slide the label sits at the
    # right end of the line on two lines; on the tax slide it sits under
    # the line on one.
    if with_tax:
        _add_text(slide, Inches(4.66), Inches(5.71), Inches(3.60),
                  Inches(0.34), "External marginal cost (EMC)", size=16,
                  bold=True, color=EMC_PURPLE, font="Calibri")
    else:
        _add_text(slide, Inches(8.22), Inches(5.40), Inches(1.66),
                  Inches(0.54), "External marginal \ncost (EMC)", size=16,
                  bold=True, color=EMC_PURPLE, font="Calibri")
    # hand-set 2026-08-30 (Nico): both curve labels at the boxes he placed
    sx, sy, sw = ((7.68, 2.59, 4.35) if with_tax else (7.53, 2.41, 2.13))
    _add_text(slide, Inches(sx), Inches(sy), Inches(sw), Inches(0.62),
              "Social marginal cost\nSMC = MC_I + EMC", size=17, bold=True,
              color=GOLD, font="Calibri")
    mx, my, mw = ((8.12, 3.32, 2.06) if with_tax else (8.12, 3.30, 2.13))
    _add_text(slide, Inches(mx), Inches(my), Inches(mw), Inches(0.57),
              "Supply = internal\nmarginal cost (MC_I)", size=17, bold=True,
              color=NAVY, font="Calibri")
    if with_tax:
        # 2026-08-30 (Nico): with a tax equal to EMC, consumers pay the
        # SMC-and-demand price while producers keep that price LESS the
        # tax - the third, lowest price line, which was missing.
        p_prod = _gas_d(GAS_Q_EXT) - 1.5
        _fig_line(slide, fig, (0, p_prod), (GAS_Q_EXT, p_prod),
                  color=GRAY, weight_pt=1.4, dash="dash")
        # P with a subscript E and a superscript P / C, as MW sets them
        _fig_ylab_subsup(slide, fig, p_prod, "P", "E", "P", size=16,
                         color=RED)
        _fig_ylab_subsup(slide, fig, _gas_d(GAS_Q_EXT), "P", "E", "C",
                         size=16, color=RED)
    # the EMC wedge, drawn where the two supply curves are furthest apart
    qw = 2.2
    _add_arrow(slide, (fig.x(qw), fig.y(_gas_mci(qw))),
               (fig.x(qw), fig.y(_gas_smc(qw))), color=EMC_PURPLE,
               weight_pt=2.5, head=True)
    if with_tax:
        _add_text(slide, Inches(5.21), Inches(2.99), Inches(3.00),
                  Inches(0.34), "Tax = EMC = $1.50", size=17, bold=True,
                  color=EMC_PURPLE, font="Calibri")
    else:
        _add_text(slide, Inches(5.75), Inches(2.79), Inches(1.29),
                  Inches(0.57), "EMC = $1.50 \nper gallon", size=17,
                  bold=True, color=EMC_PURPLE, font="Calibri")
    return fig


def slide_79_gas_wedge(prs, page_num):
    """Adopted from MW slide 68: the diagnosis on its own slide, before the
    cure.  My original carried both on one slide."""

    def draw(slide):
        _draw_gas_fig(slide, with_tax=False)
        _add_convention_box(
            slide, Inches(8.75), Inches(4.35), Inches(4.15), Inches(1.85),
            body="Drivers weigh only their own cost, so the market settles "
                 "at Q_Market. Counting the cost borne by everyone else, "
                 "the efficient quantity is the smaller Q_Ext.",
            corner_pct=0.10, size=18)

    slide = make_diagram_slide(prs, page_num, TAG_EXT,
                               "The Externality: Gasoline", draw)
    return slide


def slide_80_carbon_tax(prs, page_num):
    def draw(slide):
        _draw_gas_fig(slide, with_tax=True)
        _add_media_image(slide, "NV_s74_g13.g2.1_862b02c9.jpg",
                         left=Inches(11.25), top=Inches(1.60),
                         width=Inches(1.55))
        _add_text(slide, Inches(10.95), Inches(4.05), Inches(2.15),
                  Inches(0.30), "Arthur Pigou", size=13, italic=True,
                  color=GRAY, font="Calibri", align=PP_ALIGN.CENTER)
        # 2026-08-30 (Nico): MW slide 69's legend, verbatim — P with a
        # subscript E and a superscript C / P, then the plain-language
        # gloss, both in the same red as the ticks they explain
        for i, (sup, gloss) in enumerate(
                (("C", " : Price paid by consumer with tax on externality"),
                 ("P", " : Price received by producer with tax on "
                       "externality"))):
            # in the left margin under the plot, at 13 pt so the line
            # ends before the Q_Ext tick label at x 5.12
            box = _add_text(slide, Inches(0.55), Inches(6.52 + i * 0.30),
                            Inches(5.60), Inches(0.30), "", size=13,
                            color=RED, font="Calibri")
            p = box.text_frame.paragraphs[0]
            for txt, bl in (("P", None), ("E", "-25000"), (sup, "30000"),
                            (gloss, None)):
                r = p.add_run()
                r.text = txt
                r.font.name = "Calibri"
                r.font.size = Pt(13)
                r.font.italic = txt in ("P", "E", sup)
                r.font.color.rgb = RED
                if bl:
                    r._r.find(qn("a:rPr")).set("baseline", bl)
        _add_convention_box(
            slide, Inches(8.75), Inches(4.65), Inches(4.15), Inches(1.55),
            body="A tax equal to the external cost makes each driver face "
                 "the social marginal cost, and the market lands on Q_Ext "
                 "on its own.",
            corner_pct=0.10, size=18)

    slide = make_diagram_slide(
        prs, page_num, TAG_EXT,
        "A Carbon Tax to Correct the Negative Externality", draw)
    _set_notes(slide, NOTES[74])
    return slide


def slide_81_airplane_noise(prs, page_num):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_EXT)
    _draw_action_title(slide, "Externality: Airplane Noise")
    # 2026-08-30 (Nico): the map and the clipping get the deck's rounded
    # corners rather than sitting square
    _add_media_image(slide, "NV_s75_3_4d362b83.wmf",
                     left=Inches(0.60), top=Inches(1.60), width=Inches(7.60),
                     rounded=True, shadow=True)
    # 2026-08-28: Nico replaced the screenshot that was here with this
    # Santa Monica Lookout article, at exactly this position and size
    # (hand-edit ported from the canonical deck).
    _add_media_image(slide, "nv_smo_lookout_article.png",
                     left=Inches(8.21), top=Inches(4.10),
                     width=Inches(5.13), height=Inches(3.06),
                     rounded=True, shadow=True)
    # the caption sits right under the map so the grouping pass pairs the
    # two and they reveal on one click (course rule: a picture and its
    # caption are ONE object)
    _add_text(slide, Inches(0.60), Inches(5.66), Inches(7.60), Inches(0.30),
              "Noise contours around Santa Monica Airport", size=12,
              italic=True, color=GRAY, font="Calibri",
              align=PP_ALIGN.CENTER)
    _draw_footer(slide, FOOTER_TEXT, page_num)
    return slide


def slide_82_smo_corner(prs, page_num):
    """The corner solution: once the noise cost is added, social marginal
    cost lies above willingness to pay at every quantity, so Q* = 0."""

    def draw(slide):
        fig = SimpleFig(1.85, 6.15, 6.30, 4.10, xmax=9.4, ymax=9.0)
        _fig_axes(slide, fig, x_title="Flights per day", y_title="Price",
                  label_size=17)
        mci, smc = 2.0, 6.0
        _fig_line(slide, fig, (0, 5.0), (8.6, 0.5), color=RED,
                  weight_pt=2.75)
        # at the demand curve's own right-hand end, clear of the pilots'
        # cost label above it
        _add_text(slide, Inches(7.68), Inches(5.80),
                  Inches(3.60), Inches(0.34), "D  (MB of flying)", size=17,
                  bold=True, color=RED, font="Calibri")
        _fig_line(slide, fig, (0, mci), (8.6, mci), color=NAVY,
                  weight_pt=2.5)
        # 2026-08-30 (Nico): my original slide 76's wording — no
        # "executives", and MC with a subscript I
        _add_text(slide, Inches(7.75), Inches(4.90),
                  Inches(4.30), Inches(0.62),
                  "Pilots' internal\nMC_I of flying", size=17,
                  bold=True, color=NAVY, font="Calibri")
        _fig_line(slide, fig, (0, smc), (8.6, smc), color=GOLD,
                  weight_pt=2.75)
        _add_text(slide, Inches(7.79), Inches(3.15),
                  Inches(4.30), Inches(0.62),
                  "MC_I + tax = SMC\n(social marginal cost)", size=17,
                  bold=True, color=GOLD, font="Calibri")
        # the externality arrow takes the EMC purple of slides 81 / 82,
        # with my original's two-line label above it
        _add_arrow(slide, (fig.x(2.4), fig.y(mci)), (fig.x(2.4), fig.y(smc)),
                   color=EMC_PURPLE, weight_pt=2.5, head=True)
        _add_text(slide, Inches(3.56), Inches(3.88),
                  Inches(3.20), Inches(0.34), "Noise pollution (EMC)",
                  size=17, bold=True, color=EMC_PURPLE, font="Calibri")
        _add_text(slide, Inches(3.56), Inches(4.16),
                  Inches(3.00), Inches(0.34), "Noise tax = EMC", size=17,
                  bold=True, color=EMC_PURPLE, font="Calibri")
        # my original's Q_current: the flights flown before the tax, where
        # demand meets the pilots' own marginal cost
        q_cur = (5.0 - mci) / (4.5 / 8.6)
        _fig_line(slide, fig, (q_cur, 0), (q_cur, mci), color=NAVY,
                  weight_pt=1.4, dash="dash")
        _xlab_n(slide, fig, q_cur, "Q_current", size=16)
        _fig_dot(slide, fig, q_cur, mci, d=Inches(0.13))
        _fig_xlab(slide, fig, 0.0, "Q* = 0", size=17, bold=True)
        _fig_dot(slide, fig, 0.0, 5.0, d=Inches(0.15))
        # and the circled corner solution, top left, as in my original
        ring = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, int(fig.x(0.0) - Inches(0.22)),
            int(fig.y(5.0) - Inches(0.22)), int(Inches(0.44)),
            int(Inches(0.44)))
        ring.fill.background()
        ring.line.color.rgb = EMC_PURPLE
        ring.line.width = Pt(2.5)
        ring.shadow.inherit = False
        _add_text(slide, Inches(2.03), Inches(2.48), Inches(2.97),
                  Inches(0.57),
                  "\u201cCorner solution\u201d with Q* = 0\n"
                  "\u2192 shut down the airport", size=17, bold=True,
                  color=EMC_PURPLE, font="Calibri")
        _add_arrow(slide, (Inches(2.10), Inches(3.00)),
                   (fig.x(0.0) + Inches(0.16), fig.y(5.0) - Inches(0.20)),
                   color=EMC_PURPLE, weight_pt=1.75, head=True)

    slide = make_diagram_slide(
        prs, page_num, TAG_EXT,
        "A Pigouvian Tax for Noise? The Case of Santa Monica Airport", draw)
    _add_takeaway_bar(
        slide, "SMC lies above willingness to pay at every quantity — a "
               "corner solution at Q* = 0, so close the airport",
        top=Inches(6.58), width=Inches(11.4), height=Inches(0.52),
        left=(SLIDE_W - Inches(11.4)) // 2, fill=GOLD, text_color=NAVY,
        size=18, bold=True, rounded=True, shadow=True)
    return slide


def slide_83_summary(prs, page_num):
    bullets = [
        ("Perfect competition: the firm is a price taker, so MR = P", 0),
        ("Profit maximization: MC = MR = P gives Q*", 1),
        ("The decision to produce in the short run", 0),
        ("Continue to produce if P ≥ AVC; stop production if P < AVC", 1),
        ("The decision to stay in the industry in the long run", 0),
        ("Stay if P ≥ LAC; exit or never enter if P < LAC", 1),
        ("Long-run equilibrium: P = LMC = LAC, so profit is zero", 1),
        ("Market distortions: price controls and taxes", 0),
        ("Winners, losers, and deadweight loss", 1),
        ("Externalities, and the Pigouvian tax that corrects them", 0),
    ]
    return content_slide(prs, page_num, TAG_SUMMARY, "Module 4: Summary",
                         bullets, size=23, sub_size=21,
                         line_spacing_pts=11)


# ==========================================================================
#  ASSEMBLY
# ==========================================================================

def _poll(prs, page_num, tag, title):
    """A marked placeholder for a PollEverywhere slide.

    The title carries POLL_MARK so the pending slides stand out when
    paging through the deck; the body says what has to be done.
    """
    return make_stub(prs, page_num, tag,
                     "%s  %s" % (title, POLL_MARK), STUB_POLL)


def _strip_unused_layouts(prs):
    """Drop the python-pptx template's unused slide layouts.

    Course rule: ONE slide master for the deck, with the template defaults
    stripped.  Every slide here is built on the blank layout, so the other
    ten are dead weight that invite drift the moment someone inserts a
    slide by hand.
    """
    used = {s.slide_layout.part.partname for s in prs.slides}
    dropped = 0
    for master in prs.slide_masters:
        id_lst = master.slide_layouts._sldLayoutIdLst
        for sldLayoutId in list(id_lst):
            rId = sldLayoutId.rId
            part = master.part.related_part(rId)
            if part.partname in used:
                continue
            id_lst.remove(sldLayoutId)
            master.part.drop_rel(rId)
            dropped += 1
    return dropped


# ==========================================================================
#  SPEAKER NOTES written for this rebuild (2026-08-30)
# ==========================================================================
#  Keyed by DISPLAY number.  My original notes, carried in NOTES[...], are
#  left verbatim; these fill slides that had none and extend the thin ones.
#  Applied by _apply_written_notes() at the end of build(), which never
#  overwrites a note longer than what it would write.
WRITTEN_NOTES = {
    3: "A few housekeeping points before we start. The problem set for "
       "this module is on BruinLearn, and the practice videos are there "
       "too. Everything we do today builds on the cost curves from "
       "Module 3, so keep those handy.",
    13: "This is the intuition behind price taking. A single wheat farmer "
        "who tries to charge more than the going price simply sells "
        "nothing, because buyers have thousands of identical alternatives. "
        "That is what we mean when we say the firm faces a horizontal "
        "demand curve at the market price.",
    16: "For a price taker, revenue is simple: every extra unit sells at "
        "the same market price, so marginal revenue equals that price and "
        "never falls. This is the one case where MR = P, and it is why the "
        "profit-maximising rule collapses to P = MC. Keep that in mind, "
        "because it will not hold once we get to firms with market power.",
    18: "Here is profit maximisation seen through totals rather than "
        "margins. Total revenue is a straight line because price is "
        "constant; total cost rises and then curves up as diminishing "
        "returns bite. Profit is the vertical gap between them, and it is "
        "widest exactly where the two curves are parallel - which is the "
        "same thing as saying marginal revenue equals marginal cost.",
    19: "Meet the Yi family. They have moved from California to rural "
        "Arkansas to grow Korean cabbage and sell it to vendors in Dallas "
        "at whatever the market price happens to be. We will follow their "
        "cabbage crop for the rest of this section, so every abstract rule "
        "gets tested against a real set of numbers.",
    20: "Why this matters to you as a manager: almost no business is a "
        "textbook price taker, but many face something close to it in at "
        "least one market - a commodity input, a contract-manufacturing "
        "line, a spot market. The discipline of asking what you can and "
        "cannot influence about your price is the useful part.",
    22: "This is the Yi family's cost function. The fixed cost of 60,000 "
        "dollars is the land and equipment they are committed to for the "
        "season. The variable part rises more than proportionally with "
        "output, which is exactly the diminishing returns we saw in Module "
        "3 showing up in dollars.",
    25: "Work through the algebra with me. Marginal cost is the derivative "
        "of total cost, we set it equal to the market price, and we solve "
        "for quantity. The answer is the only output level at which the "
        "family cannot do better by producing one more or one less ton.",
    26: "Now the profit itself. We take the optimal quantity, compute "
        "total revenue and total cost at that quantity, and take the "
        "difference. Notice that the answer is positive here - that will "
        "not always be the case, and the rest of this section is about "
        "what to do when it is not.",
    27: "Average total cost is the second lens on the same problem. Total "
        "cost divided by quantity tells you what each ton costs on "
        "average, and comparing that to the price tells you immediately "
        "whether you are making money. It is the version of the "
        "calculation that most managers actually carry in their head.",
    28: "This is the picture behind the arithmetic. The height of the "
        "shaded box is price minus average total cost - the profit on each "
        "ton - and the width is the number of tons. Multiply the two and "
        "you have total profit, which is why the box is the profit.",
    31: "Here is the question that trips people up. If profit is negative, "
        "should the firm stop? Not necessarily - the fixed costs are owed "
        "either way in the short run, so the only question is whether "
        "revenue covers the costs that producing actually adds. That is "
        "the distinction the next slide makes precise.",
    32: "The short-run rule in one line. Compare price to average VARIABLE "
        "cost, not to average total cost. If price covers the variable "
        "cost of each unit, producing contributes something towards the "
        "fixed costs and the firm should keep going; if it does not, "
        "producing makes the loss bigger than the fixed costs alone.",
    33: "Now the market moves against the Yi family. Cheap Chinese imports "
        "push the US price down from 400 to 210 dollars a ton. Work "
        "through the three questions with the person next to you before we "
        "look at the answer - the numbers are all on the earlier slides.",
    37: "The new optimum. Same rule, new price: set marginal cost equal to "
        "210 and solve. The quantity falls, and when we compute profit at "
        "that quantity it is now negative - so the interesting question is "
        "the one on the next slide.",
    38: "And here is the decision. At the new price the family loses "
        "money, but price still exceeds average variable cost, so every "
        "ton they sell contributes something towards the land and "
        "equipment they are paying for anyway. They should keep producing "
        "through the season and reconsider in the long run.",
    40: "The same logic on a more realistic cost function, where average "
        "variable cost is U-shaped rather than rising throughout. Optimal "
        "output is still where marginal cost meets price. Here price sits "
        "above average total cost at that quantity, so the shaded "
        "rectangle is genuine economic profit.",
    41: "The general case at a lower price. Q* is still where MC meets P, "
        "but now price is below average total cost, so the rectangle is a "
        "loss. The firm keeps producing anyway, because price is still "
        "above average variable cost.",
    42: "The same loss case on the complex cost function. Notice that "
        "nothing about the rule changed - only where the price line "
        "happens to sit relative to the two average-cost curves.",
    43: "The lowest price case. Here the price line lies below average "
        "variable cost at every output, so there is no quantity at which "
        "producing helps. The firm stops production and eats the fixed "
        "costs.",
    44: "And the same conclusion on the complex cost function. Because "
        "marginal cost really does dip below average variable cost here, "
        "there IS a quantity where MC equals P - but price is still below "
        "AVC at that quantity, so producing would lose more than the fixed "
        "costs. Stopping is still right.",
    45: "Your turn. A coffee bean producer, marginal cost of 10, average "
        "variable cost of 12, and a market price you can read off the "
        "slide. Decide what they should do before we work it through.",
    47: "The answer. The output level is already profit-maximising, "
        "because for a price taker marginal cost equals price - so price "
        "is 10. But average variable cost is 12, which means every unit "
        "sold loses money on its own variable cost. Stop production; the "
        "fixed costs still have to be paid, and that is the smaller loss.",
    48: "Everything from this section on one slide. Find the quantity "
        "where marginal cost equals price, then use the price-versus-AVC "
        "test to decide whether to produce at all. Problem Set 3 works "
        "through both steps on new numbers.",
    52: "Now change the firm's costs and watch what happens. High-yield "
        "seeds cut the marginal cost of every ton by the same amount, so "
        "the MC curve shifts down in parallel. At the unchanged market "
        "price the farmer's optimum moves out from q0 to q1 - lower costs "
        "mean more output, not just more profit per ton.",
    59: "Now the toolbox goes somewhere less comfortable. Treat individual "
        "drug dealers as price takers in a competitive market, and the "
        "entry-and-exit logic we just built tells you what happens when "
        "policy targets sellers rather than buyers. The two slides that "
        "follow work through each case.",
    61: "Arresting dealers raises the cost of supplying, which shifts "
        "supply back. Quantity falls, but the price rises - so revenue in "
        "the market can easily go UP, and the surviving dealers make more "
        "per unit. That is the uncomfortable arithmetic behind "
        "supply-side enforcement.",
    62: "Going after users works on the other side of the market. Demand "
        "shifts back, so both quantity and price fall. Dealers make less, "
        "and the trade shrinks on both counts - which is why economists "
        "generally favour demand-side policy here.",
    63: "The long-run rules in one place. Optimal output is still where "
        "marginal cost meets price, but now the comparison is with LONG-RUN "
        "average cost, because capital is adjustable. If price cannot "
        "cover LAC the firm exits, and entry and exit push the market to "
        "the point where economic profit is zero.",
    65: "This section is about what happens when something pushes a market "
        "away from its competitive equilibrium - price controls, taxes, "
        "subsidies. The tools we need are consumer surplus, producer "
        "surplus and deadweight loss, and the next three slides build "
        "them.",
    68: "Deadweight loss is the value that simply disappears - trades that "
        "both sides would have been happy to make, which the intervention "
        "prevents. It is the difference between total welfare in the free "
        "market and total welfare under the regulation, where welfare "
        "means consumer surplus plus producer surplus.",
    70: "Here is the point that surprises most people. It does not matter "
        "which side of the market physically hands the money to the "
        "government - the burden splits the same way either way. What "
        "determines the split is elasticity: whichever side responds less "
        "to price ends up paying more of the tax.",
    71: "Before the theory, some context. Minimum wages are one of the "
        "most widely used price floors in the world, and the map shows how "
        "much the level varies across US states. Keep that variation in "
        "mind - it is what makes the empirical evidence possible.",
    73: "The minimum wage as a price floor. Set above the market wage, it "
        "raises the wage for those who keep their jobs (area A), but firms "
        "hire fewer workers, so some who would have worked do not (area "
        "C). The gap between labour supplied and labour demanded at the "
        "floor is unemployment, and B plus C is the deadweight loss.",
    74: "Before we do rent control, a quick calibration exercise. Look at "
        "the flat, the location and the map, and write down what you think "
        "it rents for. Hold your number - the answer on the next slide is "
        "usually a surprise.",
    77: "Rent control is the mirror image: a ceiling below the market "
        "price. Renters lucky enough to hold a place gain (area A), but "
        "landlords list fewer flats, so the quantity supplied falls and a "
        "shortage opens up between what renters want and what is offered. "
        "B plus C is again the deadweight loss.",
    78: "Price is not the only thing landlords can adjust. When they "
        "cannot raise the rent they can let maintenance slide, buy tenants "
        "out, or - illegally - make life unpleasant enough that people "
        "leave. The photographs are from rent-controlled buildings in "
        "Santa Monica.",
    80: "Cambridge, Massachusetts gives us an unusually clean test. Rent "
        "control ended there suddenly in 1995, so we can compare what "
        "happened to controlled units against never-controlled units "
        "nearby. Think about your prediction for each of the three "
        "questions before we look.",
    84: "The definitions we need. An externality is a cost or a benefit "
        "that lands on someone who is not part of the transaction - "
        "pollution from a power plant on one side, a well-kept front yard "
        "on the other. The gasoline number on this slide, about $1.50 a "
        "gallon in external cost, is the one we use in the worked example.",
    85: "Two families of fix. Price mechanisms - a tax on the harm or a "
        "subsidy for the benefit - change what the decision-maker faces "
        "and let the market find the quantity. Quantity mechanisms - "
        "quotas, standards, outright bans - set the quantity directly. "
        "Each has its place, and the next slides work through the price "
        "route.",
    86: "Here is the wedge. Drivers weigh only their own cost of a gallon, "
        "so the market settles where demand meets private marginal cost. "
        "Add the external cost that everyone else bears and the socially "
        "efficient quantity is smaller. The gap between the two supply "
        "curves is the external marginal cost - about $1.50 a gallon.",
    88: "Airport noise is the same problem in a form you can hear. Every "
        "flight in and out of Santa Monica imposes a cost on the "
        "neighbourhood that the pilot does not pay, and the article shows "
        "how sharply the community has pushed back.",
    89: "And here is where the arithmetic gets uncomfortable. Add the "
        "noise cost to the pilots' own cost and social marginal cost lies "
        "above willingness to pay at EVERY quantity - so the efficient "
        "number of flights is zero. This is a corner solution, and in "
        "Santa Monica it is close to what actually happened.",
    90: "Pulling it together. In the short run a price taker produces "
        "where MC equals price and keeps going as long as price covers "
        "average variable cost. In the long run entry and exit compete "
        "profit away. And once we intervene - or once a cost lands on "
        "someone outside the transaction - surplus and deadweight loss "
        "are how we judge who gains and who loses.",
    58: "Cecile Steele of Ocean View, Delaware ordered 50 chicks in 1923 "
        "and received 500 by mistake. She kept them, made a large profit, "
        "and within five years hundreds of neighbouring farmers were "
        "raising chickens for meat. That is entry in action - and it is "
        "why her extraordinary margins did not last.",
    72: "This is the podcast slide. The theory says a minimum wage should "
        "reduce employment, but the evidence is genuinely mixed: "
        "employment itself often does not fall much, while hours worked "
        "do. Listen for the distinction between experienced workers, whose "
        "total income can still rise, and low-skilled or new entrants, who "
        "tend to lose on both counts.",
    76: "Here is the answer. The unit is rent-controlled, and the "
        "Maximum Allowable Rent on the city's own database is far below "
        "what the neighbouring units fetch. That gap between the "
        "controlled rent and the market rent is the whole subject of the "
        "next few slides.",
    81: "The results. Rents rose sharply at the decontrolled units, as you "
        "would expect - but they also rose at units nearby that were never "
        "controlled, and investment went up at both. The spillover is the "
        "surprising part: rent control was holding down the value of the "
        "whole neighbourhood, not just the controlled flats.",
}


def _apply_written_notes(prs):
    """Fill in the notes above.

    Anything already on the slide is KEPT: several of my originals hold
    nothing but a source URL, which is short but not disposable, so the
    written note goes in front of it rather than over it.
    """
    n = 0
    for disp, text in WRITTEN_NOTES.items():
        if disp > len(prs.slides._sldIdLst):
            continue
        slide = prs.slides[disp - 1]
        cur = ""
        if slide.has_notes_slide:
            cur = (slide.notes_slide.notes_text_frame.text or "").strip()
        if cur and len(cur.split()) >= len(text.split()):
            continue                      # my own note already covers it
        _set_notes(slide, (text + "\n\n" + cur).strip() if cur else text)
        n += 1
    return n


def build(out_path=None):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    n = [0]

    def nxt():
        n[0] += 1
        return n[0]

    # ---- front matter --------------------------------------------------
    # the INTRODUCTION video's card opens the deck, ahead of the title
    # slide, so a video-mode deck names the video before anything else
    _video_title_card(prs, "Introduction to Market Structures", 1); nxt()
    slide_01_title(prs); nxt()
    slide_02_logistics(prs, nxt())
    make_m4_roadmap(prs, nxt())
    make_m4_outline(prs, nxt(), descriptions=True)

    # ---- 1 · introduction to market structures -------------------------
    make_m4_outline(prs, nxt(), highlight_idx=0)
    slide_06_market_structures(prs, nxt())
    slide_07_market_power(prs, nxt())

    # ---- 2 · perfect competition ---------------------------------------
    _video_title_card(prs, "Perfect Competition", 2); nxt()
    make_m4_outline(prs, nxt(), highlight_idx=1)
    slide_09_price_taker(prs, nxt())
    slide_10_market_and_firm(prs, nxt())
    slide_11_farmers(prs, nxt())

    # ---- 2a · profit maximization in the short run ---------------------
    _video_title_card(prs, "Profit Maximization of a Price Taker", 3)
    nxt()
    make_m4_outline(prs, nxt(), highlight_idx=2)
    # 2026-08-29 (Nico): the price-taker THEORY now opens topic 2a — revenue
    # conditions, then the profit-max rule, then the visual representation —
    # ahead of the Yi-family worked example.  Revenue conditions comes before
    # the rule (his order: old 18 before old 17), and Business Relevance moves
    # up ahead of the cost table so the key questions are posed before the
    # cabbage numbers arrive.  The FUNCTION names keep their original slide
    # numbers; display order is this call order, noted per line.
    slide_18_revenue_conditions(prs, nxt())      # display 13
    slide_17_profit_max_rule(prs, nxt())         # display 14
    slide_19_tr_tc_visual(prs, nxt())            # display 15
    slide_13_yi_family(prs, nxt())               # display 16
    slide_16_business_relevance(prs, nxt())      # display 17
    slide_14_cost_table(prs, nxt())              # display 18
    slide_15_tc_chart(prs, nxt())                # display 19
    slide_20_max_profit_setup(prs, nxt())
    _poll(prs, nxt(), TAG_SR, "Poll: The Optimal Quantity")
    slide_22_qstar_solution(prs, nxt())
    # 2026-08-29 (Nico): swapped by hand in PowerPoint — the worked
    # solution comes first, then the ATC restatement of profit
    slide_24_profit_solution(prs, nxt())         # display 23
    slide_23_two_ways_atc(prs, nxt())            # display 24
    slide_25_profit_rectangle(prs, nxt())
    slide_26_ross_stores(prs, nxt())
    slide_27_general_case(prs, nxt())
    slide_28_stop_producing(prs, nxt())
    # 2026-08-29 (Nico): slide_29_two_ways_avc is deleted — its
    # second identity now lives on slide 28, and the first was
    # already on display 24.  Everything below shifts up one.
    slide_30_shutdown_rule(prs, nxt())
    slide_31_new_price(prs, nxt())
    # 2026-08-30 (Nico): THREE bare PollEverywhere placeholders in place of
    # the single spliced poll — one per question on the slide before — so he
    # can paste each live poll onto an empty, correctly-chromed canvas.
    for _ in range(3):
        _poll_placeholder(prs, nxt(), TAG_SR)
    slide_33_new_qstar(prs, nxt())
    slide_34_operate_solution(prs, nxt())
    # 2026-08-30 (Nico): each price case is shown FIRST on the simple cost
    # function and then immediately on the complex one, rather than running
    # all three simple panels and then all three complex ones.
    slide_35_high_price(prs, nxt())
    slide_37b_complex_profit(prs, nxt())
    slide_36_low_price(prs, nxt())
    slide_37c_complex_loss(prs, nxt())
    slide_37_very_low_price(prs, nxt())
    slide_37d_complex_shutdown(prs, nxt())
    slide_38_coffee(prs, nxt())
    # 2026-08-30 (Nico): a bare placeholder rather than the spliced live
    # poll, so he can paste the poll onto an empty, correctly-chromed
    # canvas -- the same treatment as the three placeholders above.
    _poll_placeholder(prs, nxt(), TAG_SR)
    slide_40_coffee_solution(prs, nxt())
    slide_41_sr_summary(prs, nxt())

    # ---- 2b · firm-level and market supply -----------------------------
    _video_title_card(prs, "Firm-Level and Market Supply", 4); nxt()
    make_m4_outline(prs, nxt(), highlight_idx=3)
    slide_43_supply_curve(prs, nxt())
    slide_44_changing_mc(prs, nxt())
    slide_45_market_dynamics(prs, nxt())

    # ---- 2c · long-run competitive equilibrium -------------------------
    _video_title_card(prs, "Long-Run Competitive Equilibrium", 5); nxt()
    make_m4_outline(prs, nxt(), highlight_idx=4)
    slide_47_long_run(prs, nxt())
    slide_48_lr_equilibrium(prs, nxt())
    slide_49_chickens(prs, nxt())
    slide_51_drug_market(prs, nxt())
    # 2026-08-30 (Nico): a bare placeholder to paste the live poll onto,
    # the same treatment as slides 31-33 and 43.
    _poll_placeholder(prs, nxt(), TAG_LR)
    slide_53_arrest_dealers(prs, nxt())
    slide_54_arrest_users(prs, nxt())
    # 2026-08-30 (Nico): the long-run summary closes the subsection, just
    # ahead of the next outline slide, rather than sitting mid-block.
    slide_50_lr_summary(prs, nxt())

    # ---- 3 · market distortions and regulations ------------------------
    make_m4_outline(prs, nxt(), highlight_idx=5)
    slide_56_distortions_intro(prs, nxt())
    slide_57_consumer_surplus(prs, nxt())
    slide_58_producer_surplus(prs, nxt())
    slide_59_deadweight_loss(prs, nxt())
    slide_60_sales_tax(prs, nxt())
    slide_61_tax_incidence(prs, nxt())
    slide_62_minwage_map(prs, nxt())
    # 2026-08-30 (Nico): the podcast comes BEFORE the price-floor diagram
    slide_64_minwage_evidence(prs, nxt())
    slide_63_min_wage(prs, nxt())
    slide_65_guess_rent(prs, nxt())
    # 2026-08-30 (Nico): a bare placeholder to paste the live poll onto
    _poll_placeholder(prs, nxt(), TAG_DIST)
    slide_67_the_rent_is(prs, nxt())
    slide_68_rent_control(prs, nxt())
    # 2026-08-30 (Nico): the Prop 33 slide and its poll are dropped
    slide_71_landlord_reaction(prs, nxt())
    slide_72_cambridge_map(prs, nxt())
    slide_73_rent_questions(prs, nxt())
    slide_74_rent_results(prs, nxt())
    slide_75_argentina(prs, nxt())

    # ---- 4 · externalities ---------------------------------------------
    make_m4_outline(prs, nxt(), highlight_idx=6)
    slide_77_externalities_def(prs, nxt())
    slide_78_correcting(prs, nxt())
    slide_79_gas_wedge(prs, nxt())
    slide_80_carbon_tax(prs, nxt())
    slide_81_airplane_noise(prs, nxt())
    slide_82_smo_corner(prs, nxt())
    slide_83_summary(prs, nxt())

    # deck-wide passes
    apply_symbol_subscripts(prs)
    _apply_written_notes(prs)
    _strip_unused_layouts(prs)

    out = Path(out_path) if out_path else OUT
    prs.save(str(out))
    print("%d slides -> %s" % (len(prs.slides._sldIdLst), out))
    return out


if __name__ == "__main__":
    import sys
    build(sys.argv[1] if len(sys.argv) > 1 else None)
