"""Member-level geometry diff between two decks.

Surfaces hand edits: every shape in both decks is decoded to its RENDERED
position in inches - group children included, through the group's
off/ext / chOff/chExt transform - and matched by (type, normalised text).
Position and size are compared to 0.01".

Normalising the text strips PowerPoint's math-italic codepoints so an OMML
run does not read as a mismatch, and collapses whitespace.

Usage:  python _diff_slides.py <deck A> <deck B>
        (A = the deck that may carry hand edits, B = a fresh build)
"""
import sys
import unicodedata
from pathlib import Path

from pptx import Presentation

EMU = 914400.0
TOL = 0.01


def norm(text):
    out = []
    for ch in text or "":
        d = unicodedata.decomposition(ch)
        if d.startswith("<font>"):
            ch = chr(int(d.split()[-1], 16))
        out.append(ch)
    return " ".join("".join(out).split())


def shape_kind(sh):
    st = str(sh.shape_type or "")
    if "PICTURE" in st:
        try:
            return "pic:" + sh.image.sha1[:8]
        except Exception:
            return "pic:?"
    if sh.has_table:
        return "table"
    if sh.has_chart:
        return "chart"
    return "sp"


def walk(shapes, out, ox=0.0, oy=0.0, sx=1.0, sy=1.0):
    for sh in shapes:
        st = str(sh.shape_type or "")
        try:
            x = ox + (sh.left / EMU) * sx
            y = oy + (sh.top / EMU) * sy
            w = (sh.width / EMU) * sx
            h = (sh.height / EMU) * sy
        except TypeError:
            continue
        if "GROUP" in st:
            # children are in the group's child coordinate space
            g = sh._element
            xf = g.find(
                "{http://schemas.openxmlformats.org/presentationml/2006/main}"
                "grpSpPr")
            try:
                a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
                x2 = xf.find(a + "xfrm")
                cho = x2.find(a + "chOff")
                che = x2.find(a + "chExt")
                ext = x2.find(a + "ext")
                off = x2.find(a + "off")
                k_x = int(ext.get("cx")) / max(int(che.get("cx")), 1)
                k_y = int(ext.get("cy")) / max(int(che.get("cy")), 1)
                nox = ox + (int(off.get("x")) / EMU) * sx \
                    - (int(cho.get("x")) / EMU) * k_x * sx
                noy = oy + (int(off.get("y")) / EMU) * sy \
                    - (int(cho.get("y")) / EMU) * k_y * sy
                walk(sh.shapes, out, nox, noy, sx * k_x, sy * k_y)
            except Exception:
                walk(sh.shapes, out, ox, oy, sx, sy)
            continue
        txt = ""
        if sh.has_text_frame:
            txt = norm(sh.text_frame.text)
        elif sh.has_table:
            txt = norm(" ".join(c.text for r in sh.table.rows
                                for c in r.cells))
        out.append((shape_kind(sh), txt, round(x, 2), round(y, 2),
                    round(w, 2), round(h, 2)))


def inventory(path):
    prs = Presentation(str(path))
    slides = []
    for s in prs.slides:
        out = []
        walk(s.shapes, out)
        slides.append(out)
    return slides


def main():
    a_path, b_path = Path(sys.argv[1]), Path(sys.argv[2])
    A, B = inventory(a_path), inventory(b_path)
    print("%s: %d slides   %s: %d slides" % (a_path.name, len(A),
                                             b_path.name, len(B)))
    if len(A) != len(B):
        print("!! slide-count difference")
    diffs = 0
    for i in range(min(len(A), len(B))):
        a, b = list(A[i]), list(B[i])
        only_a, matched = [], []
        for sh in a:
            key = (sh[0], sh[1])
            hit = next((o for o in b if (o[0], o[1]) == key), None)
            if hit is None:
                only_a.append(sh)
            else:
                b.remove(hit)
                matched.append((sh, hit))
        moved = [(s, o) for s, o in matched
                 if max(abs(s[k] - o[k]) for k in (2, 3, 4, 5)) > TOL]
        if only_a or b or moved:
            diffs += 1
            print("\n--- slide %d ---" % (i + 1))
            for sh in only_a:
                print("  ONLY IN A  %-12s [%5.2f,%5.2f %5.2fx%5.2f] %s"
                      % (sh[0], sh[2], sh[3], sh[4], sh[5], sh[1][:50]))
            for sh in b:
                print("  ONLY IN B  %-12s [%5.2f,%5.2f %5.2fx%5.2f] %s"
                      % (sh[0], sh[2], sh[3], sh[4], sh[5], sh[1][:50]))
            for s, o in moved:
                print("  MOVED      %-12s A[%5.2f,%5.2f %5.2fx%5.2f] "
                      "B[%5.2f,%5.2f %5.2fx%5.2f] %s"
                      % (s[0], s[2], s[3], s[4], s[5],
                         o[2], o[3], o[4], o[5], s[1][:40]))
    print("\nslides with differences: %d" % diffs)


if __name__ == "__main__":
    main()
