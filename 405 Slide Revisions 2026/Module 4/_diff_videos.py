"""Diff every recorded video slide against its counterpart in
`Module 4 - Revised.pptx`, so Nico's hand edits in the video decks can be
adopted back into the main deck.

Compares, per mapped pair:
  * rendered shape geometry and text  (group children decoded, as
    `_diff_slides.py` does)
  * run-level formatting              (size / bold / italic / underline /
                                       colour / typeface, per paragraph)
  * speaker notes
  * grouping structure                (which shapes sit in a group)
  * animation click structure         (per-click shape signatures)

Usage:
    python _diff_videos.py            # everything
    python _diff_videos.py geom
    python _diff_videos.py runs
    python _diff_videos.py notes
    python _diff_videos.py groups
    python _diff_videos.py anim
"""
import sys
import unicodedata
import zipfile
from pathlib import Path
import xml.etree.ElementTree as ET

from pptx import Presentation

import _video_map as VM
import _match as MT

HERE = Path(__file__).resolve().parent
import os
DECK = HERE / os.environ.get("M4_DECK", "Module 4 - Revised.pptx")

EMU = 914400.0
TOL = 0.011

A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
P_NS = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
R_NS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"


# --------------------------------------------------------------- text utils

def norm(text):
    """Collapse whitespace and fold PowerPoint's math-italic codepoints."""
    out = []
    for ch in text or "":
        d = unicodedata.decomposition(ch)
        if d.startswith("<font>"):
            ch = chr(int(d.split()[-1], 16))
        out.append(ch)
    return " ".join("".join(out).split())


# ----------------------------------------------------------- geometry walk

def shape_kind(sh):
    st = str(sh.shape_type or "")
    if "PICTURE" in st:
        try:
            return "pic:" + sh.image.sha1[:8]
        except Exception:
            return "pic:?"
    if sh.has_table:
        return "table"
    if sh.has_chart:
        return "chart"
    return "sp"


def walk(shapes, out, ox=0.0, oy=0.0, sx=1.0, sy=1.0, path=""):
    for sh in shapes:
        st = str(sh.shape_type or "")
        try:
            x = ox + (sh.left / EMU) * sx
            y = oy + (sh.top / EMU) * sy
            w = (sh.width / EMU) * sx
            h = (sh.height / EMU) * sy
        except TypeError:
            continue
        if "GROUP" in st:
            g = sh._element
            xf = g.find(P_NS + "grpSpPr")
            try:
                x2 = xf.find(A_NS + "xfrm")
                cho, che = x2.find(A_NS + "chOff"), x2.find(A_NS + "chExt")
                ext, off = x2.find(A_NS + "ext"), x2.find(A_NS + "off")
                k_x = int(ext.get("cx")) / max(int(che.get("cx")), 1)
                k_y = int(ext.get("cy")) / max(int(che.get("cy")), 1)
                nox = ox + (int(off.get("x")) / EMU) * sx \
                    - (int(cho.get("x")) / EMU) * k_x * sx
                noy = oy + (int(off.get("y")) / EMU) * sy \
                    - (int(cho.get("y")) / EMU) * k_y * sy
                walk(sh.shapes, out, nox, noy, sx * k_x, sy * k_y,
                     path + "G/")
            except Exception:
                walk(sh.shapes, out, ox, oy, sx, sy, path + "G/")
            continue
        txt = ""
        if sh.has_text_frame:
            txt = norm(sh.text_frame.text)
        elif sh.has_table:
            txt = norm(" ".join(c.text for r in sh.table.rows
                                for c in r.cells))
        out.append((shape_kind(sh), txt, round(x, 2), round(y, 2),
                    round(w, 2), round(h, 2), path))


def geom(slide):
    out = []
    walk(slide.shapes, out)
    return out


# ------------------------------------------------------------ run formatting

def runs_of(slide):
    """[(shape text key, [ (para idx, run text, fmt tuple) ... ]) ...]"""
    res = []
    for sh in slide.shapes:
        for s in ([sh] if not sh.shape_type or "GROUP" not in
                  str(sh.shape_type) else list(sh.shapes)):
            if not s.has_text_frame:
                continue
            items = []
            for pi, para in enumerate(s.text_frame.paragraphs):
                for r in para.runs:
                    f = r.font
                    col = None
                    try:
                        if f.color and f.color.type is not None:
                            col = str(f.color.rgb)
                    except Exception:
                        col = "?"
                    items.append((pi, norm(r.text),
                                  (f.size.pt if f.size else None,
                                   f.bold, f.italic, f.underline, col,
                                   f.name)))
            if items:
                res.append((norm(s.text_frame.text), items))
    return res


# ------------------------------------------------------------------- notes

def notes_of(slide):
    if not slide.has_notes_slide:
        return ""
    return norm(slide.notes_slide.notes_text_frame.text)


# ------------------------------------------------------------------ groups

def group_sig(slide):
    """For each group: the sorted texts of its children."""
    sigs = []
    for sh in slide.shapes:
        if sh.shape_type and "GROUP" in str(sh.shape_type):
            kids = []
            for k in sh.shapes:
                t = norm(k.text_frame.text) if k.has_text_frame else ""
                kids.append((shape_kind(k), t[:40]))
            sigs.append(tuple(sorted(kids)))
    return sorted(sigs)


# --------------------------------------------------------------- animation

def part_map(zf):
    """display index -> slide part name"""
    pres = ET.fromstring(zf.read("ppt/presentation.xml"))
    rels = ET.fromstring(zf.read("ppt/_rels/presentation.xml.rels"))
    rmap = {c.get("Id"): c.get("Target") for c in rels}
    out = []
    for sid in pres.find(P_NS + "sldIdLst"):
        out.append("ppt/" + rmap[sid.get(R_NS + "id")].replace("../", ""))
    return out


def shape_index(root):
    """spid -> (kind, text) signature, for shapes anywhere in the tree."""
    idx = {}

    def rec(node):
        for sp in list(node):
            tag = sp.tag
            if tag in (P_NS + "sp", P_NS + "pic", P_NS + "graphicFrame",
                       P_NS + "cxnSp", P_NS + "grpSp"):
                nv = sp.find(".//" + P_NS + "cNvPr")
                if nv is not None:
                    txts = [t.text or "" for t in sp.iter(A_NS + "t")]
                    kind = tag.rsplit("}", 1)[1]
                    idx[nv.get("id")] = (kind, norm(" ".join(txts))[:44])
                if tag == P_NS + "grpSp":
                    rec(sp)
    tree = root.find(P_NS + "cSld").find(P_NS + "spTree")
    rec(tree)
    return idx


def anim_sig(xml):
    """[(trigger, [shape signature ...]) ...] — one entry per click group."""
    root = ET.fromstring(xml)
    idx = shape_index(root)
    timing = root.find(P_NS + "timing")
    if timing is None:
        return []
    clicks = []
    for par in timing.iter(P_NS + "par"):
        ctn = par.find(P_NS + "cTn")
        if ctn is None or ctn.get("nodeType") != "clickEffect":
            continue
        # this par is one click; collect every spid beneath it
        sigs = []
        for tgt in par.iter(P_NS + "spTgt"):
            sid = tgt.get("spid")
            sigs.append(idx.get(sid, ("?", "spid " + str(sid))))
        clicks.append(tuple(sigs))
    # fall back: read mainSeq children in order
    if not clicks:
        return []
    return clicks


def anim_of_all(path):
    """display index (1-based) -> click signature list"""
    zf = zipfile.ZipFile(str(path))
    parts = part_map(zf)
    out = {}
    for i, p in enumerate(parts, 1):
        out[i] = anim_sig(zf.read(p))
    zf.close()
    return out


# -------------------------------------------------------------------- main

def main():
    which = sys.argv[1] if len(sys.argv) > 1 else "all"
    rev = Presentation(str(DECK))
    rev_slides = list(rev.slides)
    rev_anim = anim_of_all(DECK) if which in ("all", "anim") else {}

    totals = {}
    for deck_name, pairs in VM.VIDEO_DECKS:
        vpath = HERE / VM.VIDEO_DIR / deck_name
        vid = Presentation(str(vpath))
        vslides = list(vid.slides)
        vanim = anim_of_all(vpath) if which in ("all", "anim") else {}
        print("\n" + "=" * 78)
        print(deck_name)
        print("=" * 78)
        for v, r in pairs:
            vs, rs = vslides[v - 1], rev_slides[r - 1]
            head = "video %2d  ->  revised %2d" % (v, r)
            msgs = []

            if which in ("all", "geom"):
                A = MT.strip_pagenum(geom(vs))
                B = MT.strip_pagenum(geom(rs))
                pairs, only_a, b = MT.match(A, B)
                moved = MT.moved(pairs)
                for sh in only_a:
                    msgs.append("  ONLY IN VIDEO %-11s [%5.2f,%5.2f "
                                "%5.2fx%5.2f] %s"
                                % (sh[0], sh[2], sh[3], sh[4], sh[5],
                                   sh[1][:46]))
                for sh in b:
                    msgs.append("  ONLY IN DECK  %-11s [%5.2f,%5.2f "
                                "%5.2fx%5.2f] %s"
                                % (sh[0], sh[2], sh[3], sh[4], sh[5],
                                   sh[1][:46]))
                for s, o in moved:
                    msgs.append("  MOVED/RESIZED %-11s V[%5.2f,%5.2f "
                                "%5.2fx%5.2f] D[%5.2f,%5.2f %5.2fx%5.2f] %s"
                                % (s[0], s[2], s[3], s[4], s[5],
                                   o[2], o[3], o[4], o[5], s[1][:34]))

            if which in ("all", "runs"):
                RA, RB = dict(runs_of(vs)), dict(runs_of(rs))
                for k in RA:
                    if k in RB and RA[k] != RB[k]:
                        for x, y in zip(RA[k], RB[k]):
                            if x != y:
                                msgs.append("  RUN FMT  %-30s V%s  D%s"
                                            % (k[:30], x, y))
                        if len(RA[k]) != len(RB[k]):
                            msgs.append("  RUN COUNT %-30s V%d D%d"
                                        % (k[:30], len(RA[k]), len(RB[k])))

            if which in ("all", "notes"):
                na, nb = notes_of(vs), notes_of(rs)
                if na != nb:
                    msgs.append("  NOTES DIFFER  video[%d] deck[%d]"
                                % (len(na), len(nb)))
                    msgs.append("     V: " + na[:150])
                    msgs.append("     D: " + nb[:150])

            if which in ("all", "groups"):
                ga, gb = group_sig(vs), group_sig(rs)
                if ga != gb:
                    msgs.append("  GROUPS DIFFER  video=%d deck=%d"
                                % (len(ga), len(gb)))
                    for g in ga:
                        if g not in gb:
                            msgs.append("     ONLY VIDEO GROUP: %s" % (g,))
                    for g in gb:
                        if g not in ga:
                            msgs.append("     ONLY DECK  GROUP: %s" % (g,))

            if which in ("all", "anim"):
                ca, cb = vanim.get(v, []), rev_anim.get(r, [])
                if ca != cb:
                    msgs.append("  ANIM DIFFERS  video clicks=%d "
                                "deck clicks=%d" % (len(ca), len(cb)))
                    for j in range(max(len(ca), len(cb))):
                        x = ca[j] if j < len(ca) else None
                        y = cb[j] if j < len(cb) else None
                        if x != y:
                            msgs.append("     click %d" % (j + 1))
                            msgs.append("        V: %s" % (x,))
                            msgs.append("        D: %s" % (y,))

            if msgs:
                totals[(deck_name, v, r)] = len(msgs)
                print("\n" + head)
                for m in msgs:
                    print(m)

    print("\n" + "=" * 78)
    print("pairs with differences: %d of 54" % len(totals))


if __name__ == "__main__":
    main()
