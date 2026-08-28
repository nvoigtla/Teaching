"""Dump a deck's per-slide text, notes, and media inventory to Markdown.

Read-only build input helper (project convention: _source_inventory.md).
Usage:  python _dump_deck.py "<deck>.pptx" <out.md>
"""
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu

EMU_IN = 914400.0


def shape_lines(sh, indent=0):
    out = []
    pad = "  " * indent
    try:
        l = sh.left / EMU_IN if sh.left is not None else None
        t = sh.top / EMU_IN if sh.top is not None else None
        w = sh.width / EMU_IN if sh.width is not None else None
        h = sh.height / EMU_IN if sh.height is not None else None
        geo = f"[{l:.2f},{t:.2f} {w:.2f}x{h:.2f}]" if None not in (l, t, w, h) else "[?]"
    except Exception:
        geo = "[?]"
    kind = str(sh.shape_type)
    if sh.shape_type is not None and "GROUP" in kind:
        out.append(f"{pad}- GROUP {geo}")
        for sub in sh.shapes:
            out += shape_lines(sub, indent + 1)
        return out
    if sh.has_text_frame:
        txt = "\n".join(p.text for p in sh.text_frame.paragraphs).strip()
        if txt:
            txt = txt.replace("\n", " ⏎ ")
            out.append(f"{pad}- TXT {geo}: {txt}")
        else:
            out.append(f"{pad}- (empty shape) {kind} {geo}")
    elif sh.shape_type is not None and "PICTURE" in kind:
        try:
            nm = sh.image.filename or sh.image.sha1[:8]
            ext = sh.image.ext
        except Exception:
            nm, ext = "?", "?"
        out.append(f"{pad}- PIC {geo} ({ext}, {nm})")
    elif sh.has_table:
        rows = []
        for r in sh.table.rows:
            rows.append(" | ".join(c.text.replace("\n", " ") for c in r.cells))
        out.append(f"{pad}- TABLE {geo}")
        for r in rows:
            out.append(f"{pad}    {r}")
    elif sh.has_chart:
        out.append(f"{pad}- CHART {geo} type={sh.chart.chart_type}")
    else:
        out.append(f"{pad}- {kind} {geo}")
    return out


def main():
    deck = Path(sys.argv[1])
    out = Path(sys.argv[2])
    prs = Presentation(str(deck))
    L = [f"# Source inventory: {deck.name}", "",
         f"Slides: {len(prs.slides)}  ·  canvas "
         f"{prs.slide_width/EMU_IN:.2f} x {prs.slide_height/EMU_IN:.2f} in", ""]
    for i, s in enumerate(prs.slides, 1):
        hidden = ""
        el = s._element
        if el.get("show") == "0":
            hidden = "  **[HIDDEN]**"
        L.append(f"## Slide {i}{hidden}")
        for sh in s.shapes:
            L += shape_lines(sh)
        if s.has_notes_slide:
            tf = s.notes_slide.notes_text_frame
            nt = tf.text.strip() if tf is not None else ""
            if nt:
                L.append("")
                L.append("**NOTES:** " + nt.replace("\n", "\n"))
        L.append("")
    out.write_text("\n".join(L), encoding="utf-8")
    print(f"{deck.name}: {len(prs.slides)} slides -> {out}")


if __name__ == "__main__":
    main()
