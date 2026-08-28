"""Extract every picture from a deck into _source_images/ + write a manifest.

Build input helper: the extracted assets are what the rebuild re-places, per
the "keep the original image assets" rule in Teaching CLAUDE.md.
"""
import sys, hashlib
from pathlib import Path
from pptx import Presentation

EMU = 914400.0
deck = Path(sys.argv[1])
outdir = Path(sys.argv[2]); outdir.mkdir(exist_ok=True)
tag = sys.argv[3]
prs = Presentation(str(deck))
rows, seen = [], {}


def walk(shapes, sn, path=""):
    for i, sh in enumerate(shapes, 1):
        st = str(sh.shape_type or "")
        if "GROUP" in st:
            walk(sh.shapes, sn, path + f"g{i}.")
            continue
        if "PICTURE" not in st:
            continue
        img = sh.image
        h = hashlib.sha1(img.blob).hexdigest()[:8]
        name = f"{tag}_s{sn:02d}_{path}{i}_{h}.{img.ext}"
        p = outdir / name
        if h in seen:
            name = seen[h]
        else:
            p.write_bytes(img.blob)
            seen[h] = name
        rows.append((sn, name, f"{sh.left/EMU:.2f},{sh.top/EMU:.2f}",
                     f"{sh.width/EMU:.2f}x{sh.height/EMU:.2f}",
                     f"{img.size[0]}x{img.size[1]}px"))


for n, s in enumerate(prs.slides, 1):
    walk(s.shapes, n)

L = [f"# Assets manifest: {deck.name}", "",
     "| Slide | File | Pos (in) | Size (in) | Native |", "|---|---|---|---|---|"]
for r in rows:
    L.append("| " + " | ".join(str(x) for x in r) + " |")
Path(f"_assets_manifest_{tag}.md").write_text("\n".join(L), encoding="utf-8")
print(f"{len(rows)} placements, {len(seen)} unique images -> {outdir}")
