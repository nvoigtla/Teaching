# ==========================================================================
#  _m4_helpers.py — shared primitive layer for the Module 4 rebuild.
#
#  Carved out of "Module 1/_build_Module1.py" by _make_helpers.py; that
#  file in turn carries the Module 7 / Module 3 helper layer verbatim.
#  Nothing here is Module-4 specific — palette, chrome, boxes, bullets,
#  OMML math, charts, tables, figures, badges, pointers, title case, and
#  the deck-wide symbol-subscript pass.
#
#  Module-4 content lives in _build_Module4.py.
# ==========================================================================

import uuid          # used by the live slide-number field footer

import copy
import math
import re
import shutil
import zipfile
from pathlib import Path

from lxml import etree as ET
from pptx import Presentation
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn
from pptx.util import Cm, Inches, Pt


# Reuse all primitives from the template script (single source of truth).
from _build_template_samples import (
    FADED,
    FOOTER_TEXT,
    GOLD,
    GOLD_W,
    GRAY,
    MARGIN,
    MODULE_AGENDA,
    NAVY,
    RULE,
    RULE_W,
    SLIDE_H,
    SLIDE_W,
    WHITE,
    _add_bulleted_list,
    _add_rect,
    _add_text,
    _blank_slide,
    _draw_action_title,
    _set_bullet_char,
)

OUT_DIR = Path(__file__).parent


# --------------------------------------------------------------------------
# Title-case top bar – replaces the all-caps default from the template script.
# The user prefers academic-paper title case (each major word capitalized,
# small connectors like "and", "of", "for", "the" stay lowercase).
# --------------------------------------------------------------------------

def _draw_top_bar_tc(slide, section_tag):
    """Navy top bar with title-cased section tag (no uppercase forcing)."""
    bar_h = Inches(0.42)
    _add_rect(slide, 0, 0, SLIDE_W, bar_h, NAVY)
    _add_text(slide, MARGIN, 0, Inches(12), bar_h,
              section_tag, size=16, bold=True,
              color=WHITE, font="Calibri",
              anchor=MSO_ANCHOR.MIDDLE)


def _draw_footer(slide, footer_text, page_num):
    """Footer rule + gold accent + footer text/page number, with larger type
    than the template default so handout printouts remain legible."""
    _add_rect(slide, 0, Inches(7.15), SLIDE_W, Inches(0.02), RULE)
    _add_rect(slide, MARGIN, Inches(7.135), GOLD_W, Inches(0.05), GOLD)
    _add_text(slide, MARGIN, Inches(7.20), Inches(11), Inches(0.32),
              footer_text, size=12, color=GRAY)
    _add_text(slide, Inches(12.5), Inches(7.20), Inches(0.6), Inches(0.32),
              str(page_num), size=12, color=GRAY, align=PP_ALIGN.RIGHT)


# --------------------------------------------------------------------------
# Speaker-notes helper
# --------------------------------------------------------------------------

def _set_notes(slide, text):
    """Replace the slide's speaker notes with *text*."""
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.clear()
    notes_tf.text = text


# --------------------------------------------------------------------------
# Reusable shape primitives for diagrams
# --------------------------------------------------------------------------

def _add_filled_box(slide, left, top, width, height, label, *,
                    fill=NAVY, text_color=WHITE, line=None,
                    size=18, bold=True, font="Calibri"):
    """Filled rectangle with centered text."""
    left, top, width, height = int(left), int(top), int(width), int(height)
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height,
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = text_color
    return shp


def _add_outlined_box(slide, left, top, width, height, label, *,
                      line=NAVY, text_color=NAVY, fill=WHITE,
                      size=18, bold=True, line_w=1.25, font="Calibri",
                      rounded=False, shadow=False, corner_pct=0.06):
    """Outlined rectangle (white fill) with centered text.

    ``rounded=True`` switches the base shape to a rounded rectangle
    (corner-adjust = ``corner_pct``); ``shadow=True`` adds a soft drop
    shadow.  Both default off so existing flat call-sites are
    unaffected.
    """
    left, top, width, height = int(left), int(top), int(width), int(height)
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(
        shape_type, left, top, width, height,
    )
    if rounded:
        try: shp.adjustments[0] = corner_pct
        except Exception: pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    if shadow:
        _add_drop_shadow(shp)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = text_color
    return shp


def _add_convention_box(slide, left, top, width, height, *,
                          prefix=None, body=None, runs=None, anchor=None,
                          space_before_pts=None,
                          fill_rgb=None, fill_alpha=None,
                          border=None, line_w=1.0,
                          corner_pct=0.12, size=15, align=PP_ALIGN.LEFT,
                          font="Calibri", pad_h=None, pad_v=None,
                          line_spacing_pct=None):
    """Cream-fill / navy-border rounded-rect explanation callout.

    The "Convention" textbox pattern from slide 14 generalised — use it
    anywhere a slide needs a compact, visually-distinct box for a
    short conceptual explanation or notational convention.  Sits well
    below a table, beside a hero formula, or as a slide-wide footer.

    Two ways to populate the text:
      • ``prefix`` (bold) + ``body`` (regular) — simplest path; matches
        slide 14's "Convention:  <text>" pattern.
      • ``runs`` — a list of ``(text, {"bold": .., "italic": .., ...})``
        tuples for finer-grained styling (multi-line, mixed formatting).

    Style defaults follow the course-layer CLAUDE.md "Convention callout
    box" spec — cream fill, thin primary-color border, slight rounding,
    primary-color text.  Override ``fill_rgb`` / ``border`` only when you
    need a different accent.
    """
    fill = fill_rgb if fill_rgb is not None else RGBColor(0xFD, 0xF6, 0xE6)
    border = border if border is not None else NAVY

    left, top, width, height = int(left), int(top), int(width), int(height)
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    # 2026-08-30 (Nico): a card can carry the SAME wash as the region it
    # describes, so the two read as one object
    if fill_alpha is not None:
        _set_fill_alpha(shp, fill_alpha)
    shp.line.color.rgb = border
    shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    try:
        shp.adjustments[0] = corner_pct
    except Exception:
        pass

    # Inset text box so the rounded corners breathe — matches slide 14.
    pad_h = Inches(0.20) if pad_h is None else pad_h
    pad_v = Inches(0.12) if pad_v is None else pad_v
    tb = slide.shapes.add_textbox(
        left + int(pad_h), top + int(pad_v),
        width - 2 * int(pad_h), height - 2 * int(pad_v),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0); tf.margin_bottom = Inches(0)
    # `anchor` lets a caller TOP-anchor the text, so that two boxes of the
    # same height line their headings up even when the bodies differ in
    # length (Module 4 slide 28's short-run / long-run pair)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE if anchor is None else anchor

    def _style_run(r, opts):
        r.font.name = opts.get('font', font)
        r.font.size = Pt(opts.get('size', size))
        r.font.bold = opts.get('bold', False)
        r.font.italic = opts.get('italic', False)
        r.font.underline = opts.get('underline', False)
        r.font.color.rgb = opts.get('color', NAVY)

    if runs is not None:
        first = True
        for entry in runs:
            text, opts = entry if isinstance(entry, tuple) else (entry, {})
            if opts.get('newline') and not first:
                p = tf.add_paragraph()
            elif first:
                p = tf.paragraphs[0]
            else:
                # Same paragraph — append run to the most-recent paragraph.
                p = tf.paragraphs[-1]
            p.alignment = align
            # 2026-08-30 (Nico): a paragraph can carry a REAL bullet - an
            # Arial round bullet with a hanging indent, the way PowerPoint
            # writes one - rather than only plain text.  python-pptx has no
            # API for this, so the pPr is written directly.
            if opts.get('bullet'):
                pPr = p._pPr if p._pPr is not None else p._p.get_or_add_pPr()
                pPr.set('marL', str(opts.get('bullet_marL', 285750)))
                pPr.set('indent', str(opts.get('bullet_indent', -285750)))
                for tag, attrs in (
                        ('a:buFont', {'typeface': 'Arial',
                                      'panose': '020B0604020202020204',
                                      'pitchFamily': '34', 'charset': '0'}),
                        ('a:buChar', {'char': '\u2022'})):
                    el = pPr.makeelement(qn(tag), attrs)
                    pPr.append(el)
            r = p.add_run()
            r.text = text
            _style_run(r, opts)
            first = False
    else:
        p = tf.paragraphs[0]
        p.alignment = align
        if prefix:
            r1 = p.add_run(); r1.text = prefix
            _style_run(r1, {'bold': True, 'color': NAVY, 'size': size})
        if body:
            r2 = p.add_run(); r2.text = body
            _style_run(r2, {'color': NAVY, 'size': size})

    if space_before_pts is not None:
        for p_obj in tf.paragraphs:
            p_obj.space_before = Pt(space_before_pts)
    if line_spacing_pct is not None:
        for p_obj in tf.paragraphs:
            pPr = p_obj._p.get_or_add_pPr()
            for old in pPr.findall(qn('a:lnSpc')):
                pPr.remove(old)
            lnSpc = ET.Element(qn('a:lnSpc'))
            spcPct = ET.SubElement(lnSpc, qn('a:spcPct'))
            spcPct.set('val', str(int(line_spacing_pct * 1000)))
            pPr.insert(0, lnSpc)
    return shp


def _set_fill_alpha(shape, opacity_pct):
    """Make a shape's SOLID fill translucent.

    ``opacity_pct`` is opacity, not transparency: 100 = fully opaque,
    30 = mostly see-through.  OOXML carries this as an ``<a:alpha>`` child
    of the colour element inside ``<a:solidFill>``, so the shape must
    already have a solid fill.  Returns the shape, so the call can be
    chained onto a box helper.
    """
    spPr = shape._element.spPr
    sf = spPr.find(qn('a:solidFill'))
    if sf is None:
        return shape
    clr = sf.find(qn('a:srgbClr'))
    if clr is None:
        return shape
    for old in clr.findall(qn('a:alpha')):
        clr.remove(old)
    alpha = ET.SubElement(clr, qn('a:alpha'))
    alpha.set('val', str(int(round(opacity_pct * 1000))))
    return shape


def _add_rounded_filled_box(slide, left, top, width, height, label, *,
                             fill=NAVY, text_color=WHITE, line=None,
                             size=18, bold=True, font="Calibri",
                             corner_pct=0.06, shadow=True):
    """Rounded-corner filled rectangle with centered text and soft drop shadow.

    Mirrors :func:`_add_filled_box` but renders ``MSO_SHAPE.ROUNDED_RECTANGLE``
    with the corner-adjust set to ``corner_pct`` (6 % per course CLAUDE.md
    "slight rounding") and a soft drop shadow via :func:`_add_drop_shadow`.
    """
    left, top, width, height = int(left), int(top), int(width), int(height)
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
    )
    try:
        shp.adjustments[0] = corner_pct
    except Exception:
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    if shadow:
        _add_drop_shadow(shp)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = text_color
    return shp


def _add_arrow(slide, start_xy, end_xy, *, color=NAVY, weight_pt=1.5,
               head=True, dash=None, head_size='med', head_both=False):
    """Draw a line/arrow from start to end (in EMU/Inches values).

    EMU coordinates MUST be integers — PowerPoint rejects decimal values
    in <a:off>/<a:ext> and refuses to open the file. Cast to int defensively.

    ``dash`` accepts any OOXML preset-dash name (e.g., ``"dash"``,
    ``"dashDot"``, ``"sysDash"``).  Default ``None`` = solid line.

    ``head_size`` is the OOXML preset arrowhead size — one of ``'sm'``,
    ``'med'`` (default), or ``'lg'``.  Width and height are set
    together, so passing ``'lg'`` gives a noticeably larger tip while
    leaving the line weight unchanged.

    ``head_both=True`` puts a triangle on BOTH ends (the market-power
    spectrum arrow on slide 7).  Schema order inside <a:ln> is
    fill → dash → join → headEnd → tailEnd, so headEnd is written first.
    """
    sx, sy = int(start_xy[0]), int(start_xy[1])
    ex, ey = int(end_xy[0]), int(end_xy[1])
    line = slide.shapes.add_connector(1, sx, sy, ex, ey)  # 1 = STRAIGHT
    line.line.color.rgb = color
    line.line.width = Pt(weight_pt)
    ln = line.line._get_or_add_ln()
    if dash is not None:
        for old in ln.findall(qn('a:prstDash')):
            ln.remove(old)
        prst = ET.SubElement(ln, qn('a:prstDash'))
        prst.set('val', dash)
    if head_both:
        headEnd = ET.SubElement(ln, qn('a:headEnd'))
        headEnd.set('type', 'triangle')
        headEnd.set('w', head_size)
        headEnd.set('h', head_size)
    if head or head_both:
        tailEnd = ET.SubElement(ln, qn('a:tailEnd'))
        tailEnd.set('type', 'triangle')
        tailEnd.set('w', head_size)
        tailEnd.set('h', head_size)
    return line


def _add_wavy_line(slide, x_start, x_end, y_center, *,
                    amplitude=None, cycles=1.75, segments=36,
                    color=NAVY, weight_pt=1.5):
    """Horizontal sinusoidal line from x_start to x_end at y_center.

    Renders a polyline approximation of ``sin(2π · t · cycles)`` inside a
    custGeom shape so the line reads as a gentle wave rather than a
    straight connector.  ``amplitude`` is the peak-to-baseline height in
    EMU; defaults to ~0.04".  ``segments`` controls how finely the wave
    is discretised — 30+ is smooth enough that the polyline reads as a
    curve.  No arrowhead.
    """
    if amplitude is None:
        amplitude = Inches(0.04)
    L = int(x_end - x_start)
    A = int(amplitude)
    if L == 0:
        return None
    flip = "1" if L < 0 else "0"
    bbox_left = int(min(x_start, x_end))
    bbox_top = int(y_center - A)
    bbox_w = abs(L)
    bbox_h = 2 * A

    pts = []
    for i in range(segments + 1):
        t = i / segments
        lx = int(round(t * 100000))
        sin_val = math.sin(2 * math.pi * t * cycles)
        # Path coords: y=0 is top.  Centre at 50000; +1·amplitude → top
        # (0), −1·amplitude → bottom (100000).
        ly = int(round(50000 - sin_val * 50000))
        pts.append((lx, ly))

    path_segs = [f'<a:moveTo><a:pt x="{pts[0][0]}" y="{pts[0][1]}"/></a:moveTo>']
    for lx, ly in pts[1:]:
        path_segs.append(f'<a:lnTo><a:pt x="{lx}" y="{ly}"/></a:lnTo>')
    path_inner = ''.join(path_segs)

    color_hex = f'{color[0]:02X}{color[1]:02X}{color[2]:02X}'
    width_emu = int(weight_pt * 12700)

    P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    A_NS_LOCAL = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    sp_xml = (
        f'<p:sp xmlns:p="{P_NS}" xmlns:a="{A_NS_LOCAL}">'
        f'<p:nvSpPr><p:cNvPr id="0" name="WavyLine"/>'
        f'<p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr>'
        f'<a:xfrm flipH="{flip}">'
        f'<a:off x="{bbox_left}" y="{bbox_top}"/>'
        f'<a:ext cx="{bbox_w}" cy="{bbox_h}"/>'
        f'</a:xfrm>'
        f'<a:custGeom>'
        f'<a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/>'
        f'<a:rect l="0" t="0" r="0" b="0"/>'
        f'<a:pathLst>'
        f'<a:path w="100000" h="100000" fill="none">'
        f'{path_inner}'
        f'</a:path>'
        f'</a:pathLst>'
        f'</a:custGeom>'
        f'<a:noFill/>'
        f'<a:ln w="{width_emu}" cap="rnd">'
        f'<a:solidFill><a:srgbClr val="{color_hex}"/></a:solidFill>'
        f'</a:ln>'
        f'</p:spPr>'
        f'</p:sp>'
    )
    elem = ET.fromstring(sp_xml)
    slide.shapes._spTree.append(elem)
    return elem


def _add_arrow_shape(slide, left, top, width, height, *,
                     direction="right", fill=GOLD, line=None):
    """Block arrow shape (the 'we-are-here' indicator).

    direction: "right" (default), "left", "up", or "down".
    """
    left, top, width, height = int(left), int(top), int(width), int(height)
    geom_map = {
        "left": MSO_SHAPE.LEFT_ARROW,
        "right": MSO_SHAPE.RIGHT_ARROW,
        "up": MSO_SHAPE.UP_ARROW,
        "down": MSO_SHAPE.DOWN_ARROW,
    }
    geom = geom_map.get(direction, MSO_SHAPE.RIGHT_ARROW)
    shp = slide.shapes.add_shape(geom, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp



# --------------------------------------------------------------------------
# Layout 2 — Section Header / Agenda (parameterized)
# --------------------------------------------------------------------------

def make_section_agenda(prs, page_num, *, current_part_idx=None,
                        current_sub_idx=None,
                        section_tag="Module 3 · Agenda",
                        title="Agenda"):
    """Render an agenda / section-divider slide.

    ``current_part_idx`` makes that Part navy and fades the others.
    ``current_sub_idx`` (optional, only meaningful with current_part_idx)
    further restricts the navy highlight inside the current Part to a
    single sub-bullet — the other subs render in faded gray, alongside
    Part 2.  Used by intra-Part dividers (e.g., the new Short Run agenda
    on page 13, or the Long Run agenda on page 31).
    """
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, section_tag)
    _draw_action_title(slide, title)

    y = Inches(1.85)
    part_title_h = Inches(0.6)
    sub_list_h = Inches(1.35)
    block_gap = Inches(0.15)

    for idx, part in enumerate(MODULE_AGENDA):
        if current_part_idx is None:
            color = NAVY  # preview: highlight all Parts
        else:
            color = NAVY if idx == current_part_idx else FADED

        # 2026-05-19: per-sub fade.  When this is the current Part AND a
        # specific sub-index is highlighted, give every other sub the
        # FADED color so the deck shows "you are HERE" within the Part.
        if (current_part_idx is not None and idx == current_part_idx
                and current_sub_idx is not None):
            sub_colors = [
                NAVY if i == current_sub_idx else FADED
                for i in range(len(part["subs"]))
            ]
        else:
            sub_colors = None

        _add_text(slide, MARGIN, y, RULE_W, part_title_h,
                  part["title"], size=30, bold=True, color=color,
                  font="Calibri")
        y += part_title_h

        _add_bulleted_list(
            slide,
            left=MARGIN + Inches(0.4),
            top=y,
            width=RULE_W - Inches(0.4),
            height=sub_list_h,
            items=part["subs"],
            size=24, color=color, bullet_color=color,
            # 2026-05-19: sub-bullets tightened from 10 → 3 pt to match
            # the deck-wide rhythm (sub-bullets cluster under their
            # parent rather than breathing out).  Applied to both
            # Parts on every agenda divider in the deck.
            line_spacing_pts=3,
            autonum_scheme='alphaLcPeriod',
            colors=sub_colors,
        )
        y += sub_list_h + block_gap

    _draw_footer(slide, FOOTER_TEXT, page_num)
    return slide


# --------------------------------------------------------------------------
# Layout 3 — Content bulleted
# --------------------------------------------------------------------------

def make_content_bulleted(prs, page_num, section_tag, title, bullets, *,
                          size=24, sub_size=None, line_spacing_pts=18,
                          sub_line_spacing_pts=None,
                          extras=None, bullets_top=None):
    """bullets: list of (text, level) tuples OR plain strings (level=0).

    ``bullets_top`` overrides the default body-region start (Inches(1.85))
    — useful when the slide also hosts large diagrams below the bullets
    and the bullets need to lift up to avoid overlap.

    ``sub_line_spacing_pts`` overrides the legacy
    ``max(6, line_spacing_pts - 8)`` formula for sub-bullet space-before.
    """
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, section_tag)
    _draw_action_title(slide, title)

    normalized = [(b, 0) if isinstance(b, str) else b for b in bullets]

    if bullets_top is None:
        bullets_top = Inches(1.85)

    _add_hierarchical_bullets(
        slide,
        left=MARGIN,
        top=bullets_top,
        width=RULE_W,
        height=Inches(5.0),
        items=normalized,
        size=size,
        sub_size=sub_size,
        line_spacing_pts=line_spacing_pts,
        sub_line_spacing_pts=sub_line_spacing_pts,
    )

    if extras is not None:
        extras(slide)

    _draw_footer(slide, FOOTER_TEXT, page_num)
    return slide


def _inject_bullet_lst_style(tf, *, size, sub_size,
                              main_color=NAVY, sub_color=GRAY,
                              main_char='▪', sub_char='–',
                              main_space_pts=None, sub_space_pts=None):
    """Inject <a:lstStyle> into a text frame defining per-level defaults.

    Defines lvl1pPr (main bullets, lvl="0") and lvl2pPr (sub bullets,
    lvl="1") with bullet character, indent, color, font, and space-
    before defaults.  Enables PowerPoint's native Tab / Shift+Tab to
    auto-reformat a bullet when its level changes: runs / paragraphs
    that do NOT explicitly override these properties inherit them.

    main_space_pts / sub_space_pts: if set, the <a:spcBef> for each
    level.  Sub-bullets default to a tighter spacing than main bullets
    (3 pt vs ~10 pt) per the deck's pedagogical preference: sub-bullets
    visually cluster under their parent, main bullets give the eye a
    larger break.
    """
    txBody = tf._txBody
    # Remove any prior lstStyle
    for old in txBody.findall(qn('a:lstStyle')):
        txBody.remove(old)
    lstStyle = ET.Element(qn('a:lstStyle'))
    levels = [
        ('a:lvl1pPr', 342900, -342900, main_color, main_char, size, main_space_pts),
        ('a:lvl2pPr', 685800, -228600, sub_color, sub_char, sub_size, sub_space_pts),
    ]
    for tag, mar_l, indent, color, char, sz, space_pts in levels:
        lvl = ET.SubElement(lstStyle, qn(tag))
        lvl.set('marL', str(mar_l))
        lvl.set('indent', str(indent))
        # spcBef must precede bullet attributes (OOXML schema order).
        if space_pts is not None:
            spcBef = ET.SubElement(lvl, qn('a:spcBef'))
            spcPts = ET.SubElement(spcBef, qn('a:spcPts'))
            spcPts.set('val', str(int(space_pts * 100)))
        bc = ET.SubElement(lvl, qn('a:buClr'))
        sc = ET.SubElement(bc, qn('a:srgbClr'))
        sc.set('val', '{:02X}{:02X}{:02X}'.format(color[0], color[1], color[2]))
        bf = ET.SubElement(lvl, qn('a:buFont'))
        bf.set('typeface', 'Calibri')
        bch = ET.SubElement(lvl, qn('a:buChar'))
        bch.set('char', char)
        drp = ET.SubElement(lvl, qn('a:defRPr'))
        drp.set('sz', str(int(sz * 100)))
        drp.set('b', '0')
        sf = ET.SubElement(drp, qn('a:solidFill'))
        sfc = ET.SubElement(sf, qn('a:srgbClr'))
        sfc.set('val', '{:02X}{:02X}{:02X}'.format(color[0], color[1], color[2]))
        lt = ET.SubElement(drp, qn('a:latin'))
        lt.set('typeface', 'Calibri')
    # Insert lstStyle after bodyPr (proper element order in <a:txBody>)
    bodyPr = txBody.find(qn('a:bodyPr'))
    if bodyPr is not None:
        bodyPr.addnext(lstStyle)
    else:
        txBody.insert(0, lstStyle)


def _add_hierarchical_bullets(slide, left, top, width, height, items,
                              *, size=24, sub_size=None, line_spacing_pts=18,
                              sub_line_spacing_pts=None):
    """Render bullets with indent levels.

    Bullet item forms:
        (text, level)                       — simple, defaults from level
        (text, level, opts)                 — opts dict overrides
        (runs_list, level, opts)            — multi-run paragraph

    text:
        - str  → a single run with text
        - ''   → empty paragraph (visual spacer; no run)
        - list → multi-run: list of ``(run_text, run_opts)`` tuples

    Paragraph-level opts (all optional):
        bullet_style: 'main' (▪ NAVY) | 'sub' (– GRAY) | 'arrow' (no bullet,
            plain left-indent; the run text supplies the leader char like
            "→" or Wingdings) | 'none'
        mar_l, indent: bullet positioning (EMU)
        space_before_pts: spcBef in pts (overrides legacy formula)
        size, color, bold, italic: defaults applied to every run unless
            the run_opts override them.

    Run-level opts (run_opts in a runs_list tuple):
        font_name (default 'Calibri'), size, color, bold, italic,
        underline (bool), wingdings (bool — emits <a:sym typeface="Wingdings"/>
        so a private-use-area character renders as its Wingdings glyph).

    Each paragraph also receives a ``lvl="N"`` attribute when level > 0
    so PowerPoint's Tab / Shift-Tab outline navigation can find an
    explicit outline level.
    """
    if sub_size is None:
        sub_size = size - 4

    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    # 2026-08-30 (Nico): python-pptx writes <a:spAutoFit/> on a new
    # textbox, which makes PowerPoint snap the height back to the text so
    # the box cannot be dragged taller by hand.  Turning autofit off also
    # makes the MIDDLE anchor mean what the docstring says it means.
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    # 2026-05-18: inject <a:lstStyle> with per-level defaults so that
    # PowerPoint's Tab / Shift+Tab keyboard shortcuts auto-reformat the
    # bullet when its outline level changes.  Default 'main'/'sub'
    # bullets (no custom indent, level < 2) skip explicit pPr/run
    # styling — they inherit from lstStyle.  Bullets with custom indent
    # OR level >= 2 OR explicit size/color overrides still apply
    # explicit styling.
    # 2026-05-18 (later): also moved space-before into lstStyle so a
    # paragraph demoted via Tab picks up the sub-bullet's tighter spcBef
    # automatically.  Sub-bullet default: 3 pt (was 6 pt) — per user
    # preference, sub-bullets cluster tightly under their parent.
    sub_default_space = (sub_line_spacing_pts if sub_line_spacing_pts is not None
                          else 3)
    _inject_bullet_lst_style(tf, size=size, sub_size=sub_size,
                              main_space_pts=line_spacing_pts,
                              sub_space_pts=sub_default_space)

    for i, item in enumerate(items):
        if len(item) == 2:
            text, level = item
            opts = {}
        else:
            text, level, opts = item
            opts = opts or {}

        # Determine inheritance up-front so spcBef logic below can use it.
        style = opts.get('bullet_style', 'main' if level == 0 else 'sub')
        has_custom_indent = 'mar_l' in opts or 'indent' in opts
        inherits_from_lst = (style in ('main', 'sub')
                              and not has_custom_indent
                              and level < 2)

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        pPr = p._p.get_or_add_pPr()

        # Space-before.  When the paragraph inherits from lstStyle and the
        # caller has not overridden, skip the per-pPr spcBef so that Tab-
        # induced level changes pick up the new level's spcBef from
        # lstStyle.  Otherwise set explicitly.
        if i > 0:
            sp_override = opts.get('space_before_pts')
            if sp_override is not None:
                spcBef = ET.SubElement(pPr, qn('a:spcBef'))
                pts = ET.SubElement(spcBef, qn('a:spcPts'))
                pts.set('val', str(sp_override * 100))
            elif not inherits_from_lst:
                if level == 0:
                    sp = line_spacing_pts
                else:
                    sp = sub_default_space
                spcBef = ET.SubElement(pPr, qn('a:spcBef'))
                pts = ET.SubElement(spcBef, qn('a:spcPts'))
                pts.set('val', str(sp * 100))
            # else: inherit from lstStyle.

        # Outline-level attribute (enables PowerPoint Tab/Shift-Tab)
        if level > 0:
            pPr.set('lvl', str(level))

        # Bullet styling — default 'main' (lvl 0) and 'sub' (lvl 1)
        # inherit from lstStyle when no custom indent is requested.
        # 'arrow' and 'none' suppress the lstStyle bullet via <a:buNone/>.
        # Level >= 2 still uses explicit _set_bullet_char (lstStyle here
        # only defines lvl1pPr and lvl2pPr).
        if style == 'main':
            if not inherits_from_lst:
                _set_bullet_char(p, char='▪', color=NAVY,
                                  mar_l=opts.get('mar_l', 342900),
                                  indent=opts.get('indent', -342900),
                                  size_pct=100)
        elif style == 'sub':
            if not inherits_from_lst:
                default_mar = 342900 + level * 342900
                _set_bullet_char(p, char='–', color=GRAY,
                                  mar_l=opts.get('mar_l', default_mar),
                                  indent=opts.get('indent', -228600),
                                  size_pct=100)
        elif style == 'arrow':
            # Plain left-indent, no bullet glyph; user text supplies leader.
            pPr.set('marL', str(opts.get('mar_l', 457200)))
            if 'indent' in opts:
                pPr.set('indent', str(opts['indent']))
            # Explicit <a:buNone/> so the lstStyle bullet doesn't bleed
            # through.
            ET.SubElement(pPr, qn('a:buNone'))
        elif style == 'none':
            ET.SubElement(pPr, qn('a:buNone'))

        # Empty paragraph (spacer) — no run; suppress the bullet so the
        # empty line doesn't render an orphan glyph.
        if text == '':
            if inherits_from_lst:
                ET.SubElement(pPr, qn('a:buNone'))
            continue

        # Normalize text to a runs list
        if isinstance(text, str):
            runs = [(text, {})]
        else:
            runs = text  # already a list of (run_text, run_opts) tuples

        # Paragraph-level explicit overrides (None = inherit from lstStyle).
        para_size_override = opts.get('size')
        para_color_override = opts.get('color')
        para_bold_override = opts.get('bold')
        para_italic_override = opts.get('italic')

        for run_text, run_opts in runs:
            run_opts = run_opts or {}

            # Inline OMML math zone — append <a14:m><m:oMath>...</m:oMath></a14:m>
            # as a sibling of the regular <a:r> runs so the formula renders
            # in-line with surrounding prose.  ``run_text`` is interpreted
            # as raw OMML content (e.g., from `_formula_mp_ratio` or the
            # `_omml_*` builders).
            if run_opts.get('omml'):
                om_sz = run_opts.get('size')
                if om_sz is None:
                    om_sz = para_size_override
                if om_sz is None:
                    om_sz = size if level == 0 else sub_size
                om_clr = run_opts.get('color') or para_color_override
                if om_clr is None:
                    om_clr = NAVY if level == 0 else GRAY
                sz_centi = int(om_sz * 100)
                clr_hex = '{:02X}{:02X}{:02X}'.format(
                    om_clr[0], om_clr[1], om_clr[2])
                a14_xml = (
                    f'<a14:m xmlns:a14="{A14_NS}" '
                    f'xmlns:m="{M_NS}" xmlns:a="{A_NS}">'
                    f'<m:oMath>{run_text}</m:oMath>'
                    f'</a14:m>'
                )
                a14_elem = ET.fromstring(a14_xml)
                for r_elem in a14_elem.iter(qn('m:r')):
                    arPr = r_elem.find(qn('a:rPr'))
                    if arPr is None:
                        arPr = ET.Element(qn('a:rPr'))
                        r_elem.insert(0, arPr)
                    arPr.set('sz', str(sz_centi))
                    if arPr.get('lang') is None:
                        arPr.set('lang', 'en-US')
                    for sf in arPr.findall(qn('a:solidFill')):
                        arPr.remove(sf)
                    # fill must come BEFORE a:latin (schema order), else
                    # PowerPoint ignores the color
                    sf = arPr.makeelement(qn('a:solidFill'), {})
                    srgb = ET.SubElement(sf, qn('a:srgbClr'))
                    srgb.set('val', clr_hex)
                    arPr.insert(0, sf)
                p._p.append(a14_elem)
                continue

            run = p.add_run()
            run.text = run_text
            run.font.name = run_opts.get('font_name', 'Calibri')

            # Size: run-opt > para-opt > inherit (for inherits_from_lst
            # cases) or fall back to level default (for explicit-styled
            # paragraphs).
            run_sz = run_opts.get('size')
            sz = run_sz if run_sz is not None else para_size_override
            if sz is None and not inherits_from_lst:
                sz = size if level == 0 else sub_size
            if sz is not None:
                run.font.size = Pt(sz)

            run_clr = run_opts.get('color')
            clr = run_clr if run_clr is not None else para_color_override
            if clr is None and not inherits_from_lst:
                clr = NAVY if level == 0 else GRAY
            if clr is not None:
                run.font.color.rgb = clr

            # Bold: only set if explicitly requested.  Default = inherit
            # (lstStyle defRPr has b="0").
            run_b = run_opts.get('bold')
            b = run_b if run_b is not None else para_bold_override
            if b is not None:
                run.font.bold = b

            it = run_opts.get('italic', para_italic_override)
            if it is not None:
                run.font.italic = it
            if run_opts.get('underline'):
                run.font.underline = True
            if run_opts.get('wingdings'):
                # Add <a:sym typeface="Wingdings"/> so private-use-area
                # characters render as their Wingdings glyphs.
                rPr = run._r.find(qn('a:rPr'))
                if rPr is None:
                    rPr = run._r.makeelement(qn('a:rPr'), {})
                    run._r.insert(0, rPr)
                sym = ET.SubElement(rPr, qn('a:sym'))
                sym.set('typeface', 'Wingdings')

    return box


# --------------------------------------------------------------------------
# Layout-3 variant — diagram canvas (action title + free-form shapes below).
# Used for the agenda flowchart and the Big Picture diagram so they live
# inside the same visual chrome as content slides but render a diagram
# instead of bullets.
# --------------------------------------------------------------------------

def make_diagram_slide(prs, page_num, section_tag, title, draw_diagram):
    """Action-title slide with free-form diagram region below.

    draw_diagram: callable(slide) that renders the diagram in the body
    region (approximately y = 1.85 to 6.95).
    """
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, section_tag)
    _draw_action_title(slide, title)
    draw_diagram(slide)
    _draw_footer(slide, FOOTER_TEXT, page_num)
    return slide


# --------------------------------------------------------------------------
# Layout 5 — Poll slide (A./B./C./D. options + POLL pill, no QR box)
# --------------------------------------------------------------------------

def _draw_poll_pill(slide, *, position='top-right',
                     fill=NAVY, text_color=WHITE, dot_color=GOLD,
                     width=None, height=None, text_size=14,
                     shadow=False):
    """Small 'POLL' pill, right-aligned.

    position: 'top-right' (default, just below the top bar) or
              'bottom-right' (bottom band aligned with discussion break).
    width / height: override the default 1.05" × 0.42" pill size.
    text_size: "POLL" font size (pt).  Scales up when the pill is enlarged.
    dot_color: pass None to skip the leading accent dot.
    shadow: add a soft drop shadow to the pill (matches discussion-break
        chrome).  Default False to preserve the original flat top-right
        pill look.
    """
    pill_w = width if width is not None else Inches(1.05)
    pill_h = height if height is not None else Inches(0.42)
    pill_x = SLIDE_W - MARGIN - pill_w
    if position == 'bottom-right':
        # Bottom-aligned with where _add_discussion_break ends
        # (top=6.25, height=0.72 → bottom=6.97).  Pill height varies, so
        # pick top so the pill BOTTOM lands on the same 6.97 baseline.
        pill_y = Inches(6.97) - pill_h
    else:
        pill_y = Inches(0.55)

    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  pill_x, pill_y, pill_w, pill_h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    if shadow:
        _add_drop_shadow(shp)
    else:
        shp.shadow.inherit = False
    try:
        shp.adjustments[0] = 0.5
    except Exception:
        pass
    _add_text(slide, pill_x, pill_y, pill_w, pill_h,
              "POLL", size=text_size, bold=True, color=text_color, font="Calibri",
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if dot_color is not None:
        dot_size = Inches(0.16)
        dot_x = pill_x - Inches(0.16) - dot_size
        dot_y = pill_y + (pill_h - dot_size) // 2
        _add_rect(slide, dot_x, dot_y, dot_size, dot_size, dot_color)


def make_poll_slide(prs, page_num, section_tag, title, options, *,
                    instructions="Respond at PollEv.com/nvoigtlaender",
                    size=30, line_spacing_pts=22):
    """Layout 5 – Poll slide.

    options: list of strings (A./B./C./D. auto-numbering applied).
    """
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, section_tag)
    _draw_action_title(slide, title)
    _draw_poll_pill(slide)

    _add_bulleted_list(
        slide,
        left=MARGIN,
        top=Inches(2.0),
        width=RULE_W,
        height=Inches(4.4),
        items=options,
        size=size, color=NAVY, bullet_color=NAVY,
        line_spacing_pts=line_spacing_pts,
        autonum_scheme='alphaUcPeriod',
    )

    _add_text(slide, MARGIN, Inches(6.55), RULE_W, Inches(0.4),
              instructions, size=16, italic=True, color=GRAY,
              align=PP_ALIGN.RIGHT)

    _draw_footer(slide, FOOTER_TEXT, page_num)
    return slide


# --------------------------------------------------------------------------
# Slide-jump hyperlinks — make a run or shape clickable in slideshow
# mode so it advances PowerPoint to a different slide.  python-pptx
# doesn't expose this directly, so we manipulate the rPr / cNvPr XML
# and register the slide relationship via the public OPC API.
# --------------------------------------------------------------------------

def _add_slide_jump_hyperlink_run(source_slide, run, target_slide,
                                    *, lock_color=True, underline=True):
    """Make ``run`` a clickable hyperlink that advances to ``target_slide``.

    ``lock_color=True`` adds the Office hyperlinkcolor extension so the
    run's explicit ``solidFill`` color is preserved at render time —
    without it, PowerPoint repaints hyperlink text in the theme's
    ``hlink`` color (a light blue) regardless of what the rPr says.

    ``underline=True`` sets ``u="sng"`` so the hyperlink stays visibly
    underlined even when the color-lock suppresses the theme styling.
    """
    rId = source_slide.part.relate_to(target_slide.part, RT.SLIDE)
    rPr = run._r.get_or_add_rPr()
    if underline:
        rPr.set('u', 'sng')
    for hl in rPr.findall(qn('a:hlinkClick')):
        rPr.remove(hl)
    hlinkClick = ET.SubElement(rPr, qn('a:hlinkClick'))
    hlinkClick.set(qn('r:id'), rId)
    hlinkClick.set('action', 'ppaction://hlinksldjump')
    if lock_color:
        AHYP_NS = 'http://schemas.microsoft.com/office/drawing/2018/hyperlinkcolor'
        ext_xml = (
            f'<a:extLst xmlns:a="{A_NS}">'
            f'<a:ext uri="{{A12FA001-AC4F-418D-AE19-62706E023703}}">'
            f'<ahyp:hlinkClr xmlns:ahyp="{AHYP_NS}" val="tx"/>'
            f'</a:ext>'
            f'</a:extLst>'
        )
        hlinkClick.append(ET.fromstring(ext_xml))


def _add_slide_jump_hyperlink_shape(source_slide, shape, target_slide):
    """Make ``shape`` clickable so the whole shape jumps to ``target_slide``."""
    rId = source_slide.part.relate_to(target_slide.part, RT.SLIDE)
    nvSpPr = shape._element.find(qn('p:nvSpPr'))
    cNvPr = nvSpPr.find(qn('p:cNvPr')) if nvSpPr is not None else None
    if cNvPr is None:
        return
    for hl in cNvPr.findall(qn('a:hlinkClick')):
        cNvPr.remove(hl)
    hlinkClick = ET.SubElement(cNvPr, qn('a:hlinkClick'))
    hlinkClick.set(qn('r:id'), rId)
    hlinkClick.set('action', 'ppaction://hlinksldjump')


# --------------------------------------------------------------------------
# Image helpers — embed source-deck images into the new deck.
# Images are pre-extracted to _source_images/slide{N}_{rId}.{ext}.
# --------------------------------------------------------------------------

SRC_IMG_DIR = Path(__file__).parent / "_source_images"


def _add_drop_shadow(shape, *, blur="50800", dist="38100",
                      direction="2700000", alpha="45000"):
    """Add a soft drop shadow to any shape that exposes spPr (pictures,
    rectangles, rounded rects).  Default: 4pt blur, 3pt offset, 45° down-
    right, 45% opacity black.  Used deck-wide for figures/boxes."""
    try:
        spPr = shape._element.spPr
    except AttributeError:
        # Fallback for shapes whose XML element exposes spPr only via find()
        spPr = shape._element.find(qn('p:spPr'))
    if spPr is None:
        return shape
    for old in spPr.findall(qn('a:effectLst')):
        spPr.remove(old)
    effLst = ET.SubElement(spPr, qn('a:effectLst'))
    outerShdw = ET.SubElement(effLst, qn('a:outerShdw'))
    outerShdw.set('blurRad', str(blur))
    outerShdw.set('dist', str(dist))
    outerShdw.set('dir', str(direction))
    outerShdw.set('algn', 'tl')
    outerShdw.set('rotWithShape', '0')
    rgb = ET.SubElement(outerShdw, qn('a:srgbClr'))
    rgb.set('val', '000000')
    a = ET.SubElement(rgb, qn('a:alpha'))
    a.set('val', str(alpha))
    return shape


def _add_graphicframe_shadow(slide, left, top, width, height, *,
                              shadow_alpha=45000):
    """White backing rectangle with an outerShdw effect, behind a
    graphicFrame (table or chart).  graphicFrames can't host
    a:effectLst directly; this rect supplies the shadow projected
    OUTSIDE its bounds.

    Charts have transparent plot areas by default, so a coloured backing
    bleeds through and tints the figure.  Use white instead — the chart
    sees white through its own transparent areas (clean background) and
    the shadow renders only at the visible edges.  Call BEFORE adding
    the table/chart so z-order is correct.
    """
    shdw = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        int(left), int(top), int(width), int(height),
    )
    shdw.fill.solid()
    shdw.fill.fore_color.rgb = WHITE
    shdw.line.fill.background()
    shdw.shadow.inherit = False
    sp_pr = shdw._element.spPr
    # Strip any default <a:effectLst/> python-pptx may have inserted
    # (duplicate effectLst makes PowerPoint refuse to open the file).
    for old in sp_pr.findall(qn('a:effectLst')):
        sp_pr.remove(old)
    effLst = ET.SubElement(sp_pr, qn('a:effectLst'))
    outerShdw = ET.SubElement(effLst, qn('a:outerShdw'))
    outerShdw.set('blurRad', '50800')
    outerShdw.set('dist', '38100')
    outerShdw.set('dir', '2700000')
    outerShdw.set('algn', 'tl')
    outerShdw.set('rotWithShape', '0')
    rgb = ET.SubElement(outerShdw, qn('a:srgbClr'))
    rgb.set('val', '000000')
    a = ET.SubElement(rgb, qn('a:alpha'))
    a.set('val', str(int(shadow_alpha)))
    return shdw


def _add_source_image(slide, src_slide_no, rid, *, left, top, width=None,
                      height=None, shadow=True, rounded=False):
    """Place a source-deck image on the new slide.

    `shadow=True` (default) adds a soft drop shadow so figures pop off the
    background — applied deck-wide per the latest visual direction.  Set
    `shadow=False` for niche cases (transparent PNGs, screenshots that
    already include a shadow, etc.).

    `rounded=True` gives real photographs the deck-standard rounded corners
    (routes through :func:`_apply_picture_style`, which sets both the
    rounded geometry and the drop shadow).  Leave False for logos,
    screenshots, and framed images per the Pictures guideline.
    """
    candidates = list(SRC_IMG_DIR.glob(f"slide{src_slide_no}_{rid}.*"))
    if not candidates:
        return None
    img = candidates[0]
    kwargs = {"left": int(left), "top": int(top)}
    if width is not None:
        kwargs["width"] = int(width)
    if height is not None:
        kwargs["height"] = int(height)
    pic = slide.shapes.add_picture(str(img), **kwargs)
    if rounded:
        _apply_picture_style(pic, corner_pct=6)
    elif shadow:
        _add_drop_shadow(pic)
    return pic


def _apply_picture_style(pic, *, corner_pct=8,
                          shadow_blur=50800, shadow_dist=38100,
                          shadow_dir=2700000, shadow_alpha=50000):
    """Apply rounded corners + drop shadow to a picture shape.

    corner_pct: rounded-corner radius as percent of shorter side (8 ≈ subtle).
    shadow_blur/dist: EMU; defaults give a soft 4pt blur, 3pt offset.
    shadow_dir: 2700000 = 45° down-right (standard).
    shadow_alpha: 50000 = 50% opacity black shadow.
    """
    spPr = pic._element.find(qn('p:spPr'))
    if spPr is None:
        return pic
    # Replace any existing prstGeom with roundRect at corner_pct
    for old in spPr.findall(qn('a:prstGeom')):
        spPr.remove(old)
    prstGeom = ET.Element(qn('a:prstGeom'))
    prstGeom.set('prst', 'roundRect')
    avLst = ET.SubElement(prstGeom, qn('a:avLst'))
    gd = ET.SubElement(avLst, qn('a:gd'))
    gd.set('name', 'adj')
    gd.set('fmla', f'val {int(corner_pct * 1000)}')
    # prstGeom must come after a:xfrm
    xfrm = spPr.find(qn('a:xfrm'))
    if xfrm is not None:
        xfrm.addnext(prstGeom)
    else:
        spPr.insert(0, prstGeom)
    # Replace any existing effectLst with a single outer shadow
    for old in spPr.findall(qn('a:effectLst')):
        spPr.remove(old)
    effectLst = ET.SubElement(spPr, qn('a:effectLst'))
    outerShdw = ET.SubElement(effectLst, qn('a:outerShdw'))
    outerShdw.set('blurRad', str(int(shadow_blur)))
    outerShdw.set('dist', str(int(shadow_dist)))
    outerShdw.set('dir', str(int(shadow_dir)))
    outerShdw.set('algn', 'tl')
    outerShdw.set('rotWithShape', '0')
    rgb = ET.SubElement(outerShdw, qn('a:srgbClr'))
    rgb.set('val', '000000')
    alpha = ET.SubElement(rgb, qn('a:alpha'))
    alpha.set('val', str(int(shadow_alpha)))
    return pic


# --------------------------------------------------------------------------
# Illustrative callout boxes — the "major-concept" badges, "Teaching Note"
# bars, and bottom-of-slide takeaway bands that recur throughout the source
# deck.  Replicating them faithfully (in template colors) is what makes the
# slides feel like the original, not a stripped-down rewrite.
# --------------------------------------------------------------------------

def _add_takeaway_bar(slide, text, *, top=Inches(6.4), width=None,
                       height=Inches(0.55), left=None,
                       fill=GOLD, text_color=WHITE,
                       size=20, font="Calibri", bold=True,
                       rounded=False, shadow=False):
    """Bottom-of-slide takeaway band — the 'major concept' callout.

    rounded=True  → ROUNDED_RECTANGLE with ~30% corner radius.
    shadow=True   → soft drop shadow (off by default to keep the flat
                    look on the majority of slides).
    left=None     → horizontally centered. Pass an EMU/Inches value to
                    pin the bar's left edge (used for hand-positioned
                    takeaways).
    """
    if width is None:
        width = Inches(9.6)
    if left is None:
        left = (SLIDE_W - width) // 2
    if not (rounded or shadow):
        return _add_filled_box(slide, left, top, width, height, text,
                                fill=fill, text_color=text_color,
                                size=size, bold=bold, font=font)
    # Inline build for the rounded / shadowed variant
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, int(left), int(top), int(width), int(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    if rounded:
        try: shape.adjustments[0] = 0.30
        except Exception: pass
    shape.shadow.inherit = False
    if shadow:
        _add_drop_shadow(shape)
    tf = shape.text_frame
    tf.margin_left = tf.margin_right = Inches(0.1)
    tf.margin_top = tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    # A newline breaks the bar onto a real second PARAGRAPH (not a soft
    # line break), which is how Nico sets a two-line takeaway by hand.
    for i, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = text_color
    return shape


def _add_teaching_note(slide, text, *, top=Inches(6.6), width=None,
                        height=Inches(0.6), left=None,
                        rounded=False, pdf_icon=False, label_color=GOLD,
                        fill_rgb=None):
    """External-document reference card.

    Visually distinct from in-slide callouts: cream/parchment fill, navy
    DASHED border, italic navy text, with a small page-icon glyph on the
    left — clearly signals 'this links to an external Teaching Note doc'.

    Optional styling for the slide-32 (page 33) treatment:
      • ``left``       — pin to an absolute X (default: horizontally center)
      • ``rounded``    — rounded corners + soft drop shadow (default flat)
      • ``pdf_icon``   — render the page glyph as a white folded-corner
                        with bold "PDF" text, sized to nearly fill the card
                        (default: small navy folded-corner, no text)
      • ``label_color``— color of the "SEE TEACHING NOTE →" prefix
                        (default GOLD; pass NAVY for the page-33 look).
      • ``fill_rgb``   — override the cream/parchment card fill
                        (default RGBColor(0xF4, 0xF1, 0xEA)).
    """
    if width is None:
        width = Inches(8.0)
    if left is None:
        left = int((SLIDE_W - width) // 2)
    left, top, width, height = int(left), int(top), int(width), int(height)

    # Card: cream fill, dashed navy border.  ``rounded=True`` switches the
    # base shape to ROUNDED_RECTANGLE and adds a soft drop shadow.
    base_shape = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    card = slide.shapes.add_shape(base_shape, left, top, width, height)
    if rounded:
        try: card.adjustments[0] = 0.20
        except Exception: pass
    card.fill.solid()
    card.fill.fore_color.rgb = (fill_rgb if fill_rgb is not None
                                  else RGBColor(0xF4, 0xF1, 0xEA))  # parchment cream default
    card.line.color.rgb = NAVY
    card.line.width = Pt(1.25)
    card.shadow.inherit = False
    if rounded:
        _add_drop_shadow(card)
    # Dashed line style
    ln = card.line._get_or_add_ln()
    # Remove any existing prstDash
    for el in ln.findall(qn('a:prstDash')):
        ln.remove(el)
    prstDash = ET.SubElement(ln, qn('a:prstDash'))
    prstDash.set('val', 'dash')

    # Empty text on the card; we add the icon + text as separate textboxes
    tf = card.text_frame
    tf.text = ""

    # Small "page" icon on the left – a folded-corner shape.  In default
    # mode it's a small filled-navy glyph (logo-ish); in pdf_icon mode it's
    # a larger WHITE sheet with thin navy border + bold "PDF" text so it
    # reads as a generic PDF document at a glance.
    icon_size = Inches(0.5) if pdf_icon else Inches(0.4)
    icon_x = left + Inches(0.2)
    icon_y = top + (height - icon_size) // 2
    page = slide.shapes.add_shape(MSO_SHAPE.FOLDED_CORNER,
                                   int(icon_x), int(icon_y),
                                   int(icon_size), int(icon_size))
    page.fill.solid()
    if pdf_icon:
        page.fill.fore_color.rgb = WHITE
        page.line.color.rgb = NAVY
        page.line.width = Pt(0.75)
    else:
        page.fill.fore_color.rgb = NAVY
        page.line.fill.background()
    page.shadow.inherit = False
    if pdf_icon:
        # "PDF" label inside the white sheet.
        ptf = page.text_frame
        ptf.word_wrap = False
        ptf.margin_left = 0
        ptf.margin_right = 0
        ptf.margin_top = 0
        ptf.margin_bottom = 0
        ptf.vertical_anchor = MSO_ANCHOR.MIDDLE
        pp = ptf.paragraphs[0]
        pp.alignment = PP_ALIGN.CENTER
        pr = pp.add_run()
        pr.text = "PDF"
        pr.font.name = "Calibri"
        pr.font.size = Pt(11)
        pr.font.bold = True
        pr.font.color.rgb = NAVY

    # Label text (italic, navy) inside the card to the right of the icon
    txt_left = int(icon_x + icon_size + Inches(0.2))
    txt_w = int(width - (icon_x + icon_size + Inches(0.4) - left))
    label_box = slide.shapes.add_textbox(txt_left, top,
                                          txt_w, height)
    label_tf = label_box.text_frame
    label_tf.word_wrap = True
    label_tf.margin_left = 0
    label_tf.margin_right = 0
    label_tf.margin_top = 0
    label_tf.margin_bottom = 0
    label_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = label_tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    # First run: "See teaching note:" prefix in caller-supplied color
    # (GOLD by default; NAVY for the slide-33 treatment).
    r1 = p.add_run()
    r1.text = "SEE TEACHING NOTE  →  "
    r1.font.name = "Calibri"
    r1.font.size = Pt(12)
    r1.font.bold = True
    r1.font.color.rgb = label_color
    # Second run: the title of the note, italic navy
    r2 = p.add_run()
    r2.text = text
    r2.font.name = "Calibri"
    r2.font.size = Pt(16)
    r2.font.italic = True
    r2.font.bold = True
    r2.font.color.rgb = NAVY
    return card


def _add_discussion_break(slide, *, top=Inches(6.25), width=Inches(4.8),
                           left=None, text="Discussion Break"):
    """Rounded-parallelogram 'discussion break' badge (bottom-right).

    Custom-geometry shape: top and bottom edges are horizontal; the left
    and right edges slant at 45° in real space (skew = height of the
    shape).  All four corners are slightly rounded.  Gold fill, navy
    bold text, soft drop shadow.

    left=None → pin to the right edge with the default MARGIN. Pass an
    EMU/Inches value to override (used for hand-positioned badges).
    """
    height = Inches(0.72)
    if left is None:
        left = SLIDE_W - MARGIN - width
    left, top, width, height = int(left), int(top), int(width), int(height)
    # Compute skew in path-coordinate units so that left/right sides slant
    # at 45° in REAL space:  horizontal offset of top edge == shape height.
    skew = int(100000 * height / width) if width else 15000
    skew = min(max(skew, 6000), 35000)
    r = 5000          # corner radius in path units — gentle rounding

    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = GOLD
    shp.line.fill.background()
    shp.shadow.inherit = False
    # Strip the default prstGeom – we'll inject custGeom instead.
    spPr = shp._element.spPr
    for old in spPr.findall(qn('a:prstGeom')):
        spPr.remove(old)
    rs = int(r * skew / 100000)
    # IMPORTANT: <a:rect> defines the TEXT bounding rectangle inside the
    # custom geometry.  Set it to the parallelogram's inscribed rectangle
    # (from TL vertex to BR vertex) so PowerPoint won't render text past
    # the slanted edges, regardless of the text frame's own lIns/rIns.
    custgeom_xml = (
        f'<a:custGeom xmlns:a="{A_NS}">'
        f'<a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/>'
        f'<a:rect l="{skew}" t="0" r="{100000-skew}" b="100000"/>'
        f'<a:pathLst><a:path w="100000" h="100000">'
        # Top-left rounded corner (vertex at (skew, 0))
        f'<a:moveTo><a:pt x="{skew+r}" y="0"/></a:moveTo>'
        # Top edge
        f'<a:lnTo><a:pt x="{100000-r}" y="0"/></a:lnTo>'
        # Top-right corner round
        f'<a:cubicBezTo>'
        f'<a:pt x="100000" y="0"/><a:pt x="100000" y="0"/>'
        f'<a:pt x="{100000-rs}" y="{r}"/>'
        f'</a:cubicBezTo>'
        # Right slanted side (down-left)
        f'<a:lnTo><a:pt x="{100000-skew+rs}" y="{100000-r}"/></a:lnTo>'
        # Bottom-right corner round (vertex at (100000-skew, 100000))
        f'<a:cubicBezTo>'
        f'<a:pt x="{100000-skew}" y="100000"/>'
        f'<a:pt x="{100000-skew}" y="100000"/>'
        f'<a:pt x="{100000-skew-r}" y="100000"/>'
        f'</a:cubicBezTo>'
        # Bottom edge
        f'<a:lnTo><a:pt x="{r}" y="100000"/></a:lnTo>'
        # Bottom-left corner round (vertex at (0, 100000))
        f'<a:cubicBezTo>'
        f'<a:pt x="0" y="100000"/><a:pt x="0" y="100000"/>'
        f'<a:pt x="{rs}" y="{100000-r}"/>'
        f'</a:cubicBezTo>'
        # Left slanted side (up-right)
        f'<a:lnTo><a:pt x="{skew-rs}" y="{r}"/></a:lnTo>'
        # Top-left corner round (close the path)
        f'<a:cubicBezTo>'
        f'<a:pt x="{skew}" y="0"/><a:pt x="{skew}" y="0"/>'
        f'<a:pt x="{skew+r}" y="0"/>'
        f'</a:cubicBezTo>'
        f'<a:close/>'
        f'</a:path></a:pathLst>'
        f'</a:custGeom>'
    )
    custgeom = ET.fromstring(custgeom_xml)
    # Insert custGeom right after a:xfrm (schema order)
    xfrm = spPr.find(qn('a:xfrm'))
    if xfrm is not None:
        xfrm.addnext(custgeom)
    else:
        spPr.insert(0, custgeom)

    # Drop shadow (45° down-right, 50% opacity).
    for old in spPr.findall(qn('a:effectLst')):
        spPr.remove(old)
    effectLst = ET.SubElement(spPr, qn('a:effectLst'))
    outerShdw = ET.SubElement(effectLst, qn('a:outerShdw'))
    outerShdw.set('blurRad', '50800')
    outerShdw.set('dist', '38100')
    outerShdw.set('dir', '2700000')
    outerShdw.set('algn', 'tl')
    outerShdw.set('rotWithShape', '0')
    rgb = ET.SubElement(outerShdw, qn('a:srgbClr'))
    rgb.set('val', '000000')
    alpha = ET.SubElement(rgb, qn('a:alpha'))
    alpha.set('val', '50000')

    # Text — render as a SEPARATE textbox overlaid on top of the
    # parallelogram, positioned exactly inside the inscribed rectangle.
    # This decouples text placement from the shape geometry and avoids
    # PowerPoint rendering the run past the slanted edges (which can
    # happen with the in-shape text frame on some PowerPoint versions
    # even with <a:rect> set inside custGeom).
    # In real EMU, the skew equals the shape height (45° slant), so the
    # inscribed rectangle spans (left + height) → (left + width - height).
    skew_emu = height
    ins_left = left + skew_emu
    ins_top = top
    ins_w = width - 2 * skew_emu
    ins_h = height
    txt = slide.shapes.add_textbox(int(ins_left), int(ins_top),
                                     int(ins_w), int(ins_h))
    tf = txt.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(28)         # bumped 20 → 28 on 2026-05-16
    run.font.bold = True
    run.font.color.rgb = NAVY
    shp.name = "sdbadge:box"
    txt.name = "sdbadge:txt"
    return shp


def _add_callout_box(slide, left, top, width, height, text, *,
                      fill=GOLD, text_color=WHITE, size=14, bold=True,
                      rounded=False, shadow=False):
    """Small free-form annotation/callout (e.g., 'plot the slope', 'Revenue
    per car net of material cost').  Used to mark a graph or sub-region.

    ``rounded``/``shadow`` (default off) route to the rounded + drop-shadow
    variant for the deck-wide box treatment."""
    if rounded or shadow:
        return _add_rounded_filled_box(
            slide, left, top, width, height, text,
            fill=fill, text_color=text_color,
            size=size, bold=bold, font="Calibri",
            corner_pct=0.12, shadow=shadow)
    return _add_filled_box(slide, left, top, width, height, text,
                            fill=fill, text_color=text_color,
                            size=size, bold=bold, font="Calibri")


def _add_anchor_burst(slide, left, top, width, height,
                       top_text, bottom_text=None, extra_text=None,
                       *, fill=GOLD, text_color=NAVY,
                       top_size=14, bottom_size=10):
    """12-point star background + a separate text-box overlay.

    The star is purely decorative; the text lives in a normal
    rectangular text box layered on top so the text isn't constrained
    by the star's geometry (no risk of bleeding into the points).
    Reusable on every slide where MB = MC is being invoked, so the same
    visual pattern carries over.
    """
    left, top, width, height = int(left), int(top), int(width), int(height)

    # 1. Decorative star background (no text)
    shp = slide.shapes.add_shape(MSO_SHAPE.STAR_12_POINT, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = NAVY
    shp.line.width = Pt(1.0)
    shp.shadow.inherit = False
    # Soft drop shadow — added 2026-05-15 per user request, so the MB=MC
    # star reads as lifted off the slide like every other content shape.
    _add_drop_shadow(shp)
    # Suppress any auto-inserted text frame contents on the shape itself.
    shp.text_frame.text = ""

    # 2. Overlay text box, sized to the star's inscribed body
    inner_w = int(width * 0.65)
    inner_h = int(height * 0.55)
    inner_x = left + (width - inner_w) // 2
    inner_y = top + (height - inner_h) // 2

    box = slide.shapes.add_textbox(inner_x, inner_y, inner_w, inner_h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = top_text
    r1.font.name = 'Calibri'
    r1.font.size = Pt(top_size)
    r1.font.bold = True
    r1.font.color.rgb = text_color
    if bottom_text:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = bottom_text
        r2.font.name = 'Calibri'
        r2.font.size = Pt(bottom_size)
        r2.font.italic = True
        r2.font.bold = True
        r2.font.color.rgb = text_color
    if extra_text:
        p3 = tf.add_paragraph()
        p3.alignment = PP_ALIGN.CENTER
        r3 = p3.add_run()
        r3.text = extra_text
        r3.font.name = 'Calibri'
        r3.font.size = Pt(bottom_size)
        r3.font.italic = True
        r3.font.bold = True
        r3.font.color.rgb = text_color
    return shp


# --------------------------------------------------------------------------
# OMML (Office Math Markup Language) equation helper – gives formulas a
# proper TeX-style render with italic variables, stacked fractions, real
# subscripts/superscripts.  Uses Cambria Math (the standard PPT math font).
# --------------------------------------------------------------------------

M_NS = 'http://schemas.openxmlformats.org/officeDocument/2006/math'
A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
A14_NS = 'http://schemas.microsoft.com/office/drawing/2010/main'


def _omml_fill(color):
    """Build the inner ``<a:solidFill>`` clause for an OMML run's rPr.

    ``color`` may be an ``RGBColor`` instance (or a 3-tuple of ints).
    Returns an empty string when ``color`` is None so callers can splice
    the result into rPr unconditionally.
    """
    if color is None:
        return ''
    return (
        f'<a:solidFill><a:srgbClr val="'
        f'{color[0]:02X}{color[1]:02X}{color[2]:02X}'
        f'"/></a:solidFill>'
    )


def _omml_run(text, *, color=None, bold=False):
    """OMML run for an italic variable (default math style).

    Inside an oMath, italic style is the math default for Latin letters;
    we leave m:rPr out entirely so the Cambria Math italic comes through.
    The a:rPr applies drawing-level font sizing/coloring.  Pass ``color``
    to tint the run (e.g., green ΔL / ΔQ in the slide-14 Convention box).
    """
    return (
        f'<m:r xmlns:m="{M_NS}">'
        f'<a:rPr xmlns:a="{A_NS}" lang="en-US" b="{int(bold)}" i="1">'
        f'{_omml_fill(color)}'
        f'<a:latin typeface="Cambria Math"/>'
        f'<a:ea typeface="Cambria Math"/>'
        f'</a:rPr>'
        f'<m:t>{text}</m:t>'
        f'</m:r>'
    )


def _omml_text(text, *, color=None, bold=False):
    """Upright-style OMML run (for operators, numbers, acronyms).

    Force plain (upright) style via <m:rPr><m:sty m:val="p"/></m:rPr> – this
    is the documented way to disable the math-default italics for the
    enclosed run.  Pass ``color`` to tint the run.
    """
    return (
        f'<m:r xmlns:m="{M_NS}">'
        f'<m:rPr><m:sty m:val="p"/></m:rPr>'
        f'<a:rPr xmlns:a="{A_NS}" lang="en-US" b="{int(bold)}" i="0">'
        f'{_omml_fill(color)}'
        f'<a:latin typeface="Cambria Math"/>'
        f'</a:rPr>'
        f'<m:t xml:space="preserve">{text}</m:t>'
        f'</m:r>'
    )


def _omml_sub(base, sub):
    """OMML subscript: base with subscript expression."""
    return (
        f'<m:sSub xmlns:m="{M_NS}">'
        f'<m:sSubPr><m:ctrlPr>'
        f'<a:rPr xmlns:a="{A_NS}" lang="en-US" i="1">'
        f'<a:latin typeface="Cambria Math"/></a:rPr>'
        f'</m:ctrlPr></m:sSubPr>'
        f'<m:e>{base}</m:e>'
        f'<m:sub>{sub}</m:sub>'
        f'</m:sSub>'
    )


def _omml_underbrace(expr, label):
    """OMML under-brace with a label beneath it — TeX's

        \\underbrace{expr}_{label}

    This is a <m:limLow> (lower limit) wrapped around a <m:groupChr> whose
    character is U+23DF, the bottom curly bracket; that is exactly the
    structure Word writes, so PowerPoint renders and edits it natively.
    """
    return (
        f'<m:limLow xmlns:m="{M_NS}">'
        f'<m:limLowPr><m:ctrlPr>'
        f'<a:rPr xmlns:a="{A_NS}" lang="en-US" i="1">'
        f'<a:latin typeface="Cambria Math"/></a:rPr>'
        f'</m:ctrlPr></m:limLowPr>'
        f'<m:e>'
        f'<m:groupChr>'
        f'<m:groupChrPr>'
        f'<m:chr m:val="&#9183;"/>'
        f'<m:pos m:val="bot"/>'
        f'<m:vertJc m:val="top"/>'
        f'<m:ctrlPr>'
        f'<a:rPr xmlns:a="{A_NS}" lang="en-US" i="1">'
        f'<a:latin typeface="Cambria Math"/></a:rPr>'
        f'</m:ctrlPr>'
        f'</m:groupChrPr>'
        f'<m:e>{expr}</m:e>'
        f'</m:groupChr>'
        f'</m:e>'
        f'<m:lim>{label}</m:lim>'
        f'</m:limLow>'
    )


def _omml_frac(num, den):
    """OMML stacked fraction: num / den."""
    return (
        f'<m:f xmlns:m="{M_NS}">'
        f'<m:fPr><m:ctrlPr>'
        f'<a:rPr xmlns:a="{A_NS}" lang="en-US" i="1">'
        f'<a:latin typeface="Cambria Math"/></a:rPr>'
        f'</m:ctrlPr></m:fPr>'
        f'<m:num>{num}</m:num>'
        f'<m:den>{den}</m:den>'
        f'</m:f>'
    )


def _omml_sup(base, sup):
    """OMML superscript: base^sup (e.g. Q²)."""
    return (
        f'<m:sSup xmlns:m="{M_NS}">'
        f'<m:sSupPr><m:ctrlPr>'
        f'<a:rPr xmlns:a="{A_NS}" lang="en-US" i="1">'
        f'<a:latin typeface="Cambria Math"/></a:rPr>'
        f'</m:ctrlPr></m:sSupPr>'
        f'<m:e>{base}</m:e>'
        f'<m:sup>{sup}</m:sup>'
        f'</m:sSup>'
    )


def _add_dashed_gridlines(axis_el):
    """Dashed light-grey major gridlines on an axis (Cobb-Douglas style)."""
    for old in axis_el.findall(qn('c:majorGridlines')):
        axis_el.remove(old)
    gl = ET.Element(qn('c:majorGridlines'))
    sp = ET.SubElement(gl, qn('c:spPr'))
    ln = ET.SubElement(sp, qn('a:ln'))
    ln.set('w', '9525')
    ln.set('cap', 'flat'); ln.set('cmpd', 'sng'); ln.set('algn', 'ctr')
    fill = ET.SubElement(ln, qn('a:solidFill'))
    clr = ET.SubElement(fill, qn('a:srgbClr')); clr.set('val', 'C8CDD3')
    dash = ET.SubElement(ln, qn('a:prstDash')); dash.set('val', 'dash')
    axpos = axis_el.find(qn('c:axPos'))
    if axpos is not None:
        axpos.addnext(gl)


def _align_x_labels_with_ticks(value_axis):
    """Set <c:crossBetween val="midCat"/> on a value axis so the category
    labels and tick marks align (default for line charts in OOXML is
    "between", which leaves labels visually between adjacent ticks).
    """
    val_el = value_axis._element
    for old in val_el.findall(qn('c:crossBetween')):
        val_el.remove(old)
    cb = ET.Element(qn('c:crossBetween'))
    cb.set('val', 'midCat')
    # Schema position: c:crossBetween follows c:crosses(At) and precedes
    # c:majorUnit.  Insert before c:majorUnit if present; else append.
    mu_el = val_el.find(qn('c:majorUnit'))
    if mu_el is not None:
        mu_el.addprevious(cb)
    else:
        val_el.append(cb)


def _make_simple_line_chart(slide, x, y, w, h, categories, values, *,
                              line_color, x_title, y_title,
                              y_min=0, y_max=None, y_unit=None,
                              marker='circle'):
    """Single-series line+markers chart with dashed light-grey gridlines.

    Same visual conventions as slide 11 (Calibri navy labels, dashed C8CDD3
    major gridlines, marker size 7) but no legend/title.
    """
    cd = CategoryChartData()
    cd.categories = list(categories)
    cd.add_series("Y", values)
    # Drop-shadow rectangle behind the chart (graphicFrames can't host shadow).
    _add_graphicframe_shadow(slide, x, y, w, h)
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE,
        int(x), int(y), int(w), int(h), cd,
    )
    chart = chart_shape.chart
    chart.has_title = False
    chart.has_legend = False

    clr_hex = f'{line_color[0]:02X}{line_color[1]:02X}{line_color[2]:02X}'

    for series in chart.series:
        line = series.format.line
        line.color.rgb = line_color
        line.width = Pt(2.5)
        ser_xml = series._element
        for old in ser_xml.findall(qn('c:marker')):
            ser_xml.remove(old)
        m = ET.SubElement(ser_xml, qn('c:marker'))
        sym = ET.SubElement(m, qn('c:symbol')); sym.set('val', marker)
        sz_el = ET.SubElement(m, qn('c:size')); sz_el.set('val', '7')
        sp = ET.SubElement(m, qn('c:spPr'))
        fl = ET.SubElement(sp, qn('a:solidFill'))
        rg = ET.SubElement(fl, qn('a:srgbClr')); rg.set('val', clr_hex)
        ln = ET.SubElement(sp, qn('a:ln'))
        lf = ET.SubElement(ln, qn('a:solidFill'))
        lr = ET.SubElement(lf, qn('a:srgbClr')); lr.set('val', clr_hex)
        # disable smoothing (straight segments between points)
        for sm in ser_xml.findall(qn('c:smooth')):
            ser_xml.remove(sm)
        smooth = ET.SubElement(ser_xml, qn('c:smooth'))
        smooth.set('val', '0')

    # Axes – axis titles in BOLD ITALIC navy (per course CLAUDE.md);
    # tick labels in regular Calibri navy.
    cat = chart.category_axis
    cat.tick_labels.font.name = "Calibri"
    cat.tick_labels.font.size = Pt(10)
    cat.tick_labels.font.color.rgb = NAVY
    cat.has_title = True
    cat.axis_title.text_frame.text = x_title
    ar = cat.axis_title.text_frame.paragraphs[0].runs[0]
    ar.font.name = "Calibri"; ar.font.size = Pt(12)
    ar.font.bold = True; ar.font.italic = True
    ar.font.color.rgb = NAVY

    val = chart.value_axis
    val.tick_labels.font.name = "Calibri"
    val.tick_labels.font.size = Pt(10)
    val.tick_labels.font.color.rgb = NAVY
    val.has_title = True
    val.axis_title.text_frame.text = y_title
    ar = val.axis_title.text_frame.paragraphs[0].runs[0]
    ar.font.name = "Calibri"; ar.font.size = Pt(12)
    ar.font.bold = True; ar.font.italic = True
    ar.font.color.rgb = NAVY
    if y_max is not None:
        val.minimum_scale = y_min
        val.maximum_scale = y_max
    if y_unit is not None:
        val.major_unit = y_unit

    # Align X-axis category labels with tick marks (default OOXML places
    # them in the gaps between ticks).
    _align_x_labels_with_ticks(val)

    _add_dashed_gridlines(cat._element)
    _add_dashed_gridlines(val._element)
    return chart_shape


def _make_multi_line_chart(slide, x, y, w, h, categories, series, *,
                             x_title, y_title,
                             y_min=0, y_max=None, y_unit=None,
                             legend=True, legend_pos=('0.08', '0.10', '0.20', '0.20')):
    """Multi-series line+markers chart with the deck's standard styling.

    series: list of (name, values, color: RGBColor, marker: str) tuples.
    legend_pos: (x, y, w, h) in chart-fraction units (str) – top-left default.
    """
    _add_graphicframe_shadow(slide, x, y, w, h)
    cd = CategoryChartData()
    cd.categories = list(categories)
    for name, values, _color, _marker in series:
        cd.add_series(name, values)
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE,
        int(x), int(y), int(w), int(h), cd,
    )
    chart = chart_shape.chart
    chart.has_title = False
    chart.has_legend = legend

    # Per-series styling: line color + marker
    for idx, ser in enumerate(chart.series):
        name, values, color, marker = series[idx]
        line = ser.format.line
        line.color.rgb = color
        line.width = Pt(2.5)
        ser_xml = ser._element
        clr_hex = f'{color[0]:02X}{color[1]:02X}{color[2]:02X}'
        for old in ser_xml.findall(qn('c:marker')):
            ser_xml.remove(old)
        m = ET.SubElement(ser_xml, qn('c:marker'))
        sym = ET.SubElement(m, qn('c:symbol')); sym.set('val', marker)
        sz_el = ET.SubElement(m, qn('c:size')); sz_el.set('val', '7')
        sp = ET.SubElement(m, qn('c:spPr'))
        fl = ET.SubElement(sp, qn('a:solidFill'))
        rg = ET.SubElement(fl, qn('a:srgbClr')); rg.set('val', clr_hex)
        ln = ET.SubElement(sp, qn('a:ln'))
        lf = ET.SubElement(ln, qn('a:solidFill'))
        lr = ET.SubElement(lf, qn('a:srgbClr')); lr.set('val', clr_hex)
        for sm in ser_xml.findall(qn('c:smooth')):
            ser_xml.remove(sm)
        smooth = ET.SubElement(ser_xml, qn('c:smooth'))
        smooth.set('val', '0')

    # Axes – bold italic navy titles per course style
    cat = chart.category_axis
    cat.tick_labels.font.name = "Calibri"
    cat.tick_labels.font.size = Pt(10)
    cat.tick_labels.font.color.rgb = NAVY
    cat.has_title = True
    cat.axis_title.text_frame.text = x_title
    ar = cat.axis_title.text_frame.paragraphs[0].runs[0]
    ar.font.name = "Calibri"; ar.font.size = Pt(12)
    ar.font.bold = True; ar.font.italic = True; ar.font.color.rgb = NAVY

    val = chart.value_axis
    val.tick_labels.font.name = "Calibri"
    val.tick_labels.font.size = Pt(10)
    val.tick_labels.font.color.rgb = NAVY
    val.has_title = True
    val.axis_title.text_frame.text = y_title
    ar = val.axis_title.text_frame.paragraphs[0].runs[0]
    ar.font.name = "Calibri"; ar.font.size = Pt(12)
    ar.font.bold = True; ar.font.italic = True; ar.font.color.rgb = NAVY
    if y_max is not None:
        val.minimum_scale = y_min
        val.maximum_scale = y_max
    if y_unit is not None:
        val.major_unit = y_unit

    # Align X-axis category labels with tick marks.
    _align_x_labels_with_ticks(val)

    _add_dashed_gridlines(cat._element)
    _add_dashed_gridlines(val._element)

    # Legend top-left inside plot with white fill
    if legend:
        leg_el = chart.legend._element
        chart.legend.font.name = "Calibri"
        chart.legend.font.size = Pt(11)
        chart.legend.font.color.rgb = NAVY
        chart.legend.include_in_layout = False
        # Strip default legendPos / layout, replace
        for old in leg_el.findall(qn('c:layout')):
            leg_el.remove(old)
        for old in leg_el.findall(qn('c:legendPos')):
            leg_el.remove(old)
        pos = ET.SubElement(leg_el, qn('c:legendPos')); pos.set('val', 'tr')
        leg_el.remove(pos); leg_el.insert(0, pos)
        # manualLayout positions in chart-fraction units
        layout = ET.Element(qn('c:layout'))
        ml = ET.SubElement(layout, qn('c:manualLayout'))
        ET.SubElement(ml, qn('c:xMode')).set('val', 'edge')
        ET.SubElement(ml, qn('c:yMode')).set('val', 'edge')
        ET.SubElement(ml, qn('c:x')).set('val', legend_pos[0])
        ET.SubElement(ml, qn('c:y')).set('val', legend_pos[1])
        ET.SubElement(ml, qn('c:w')).set('val', legend_pos[2])
        ET.SubElement(ml, qn('c:h')).set('val', legend_pos[3])
        pos.addnext(layout)
        # White fill behind legend
        for old in leg_el.findall(qn('c:spPr')):
            leg_el.remove(old)
        leg_spPr = ET.Element(qn('c:spPr'))
        sf = ET.SubElement(leg_spPr, qn('a:solidFill'))
        clr = ET.SubElement(sf, qn('a:srgbClr')); clr.set('val', 'FFFFFF')
        ln = ET.SubElement(leg_spPr, qn('a:ln')); ln.set('w', '6350')
        lf = ET.SubElement(ln, qn('a:solidFill'))
        lc = ET.SubElement(lf, qn('a:srgbClr')); lc.set('val', '0B2B4E')
        layout.addnext(leg_spPr)
    return chart_shape


def _make_xy_line_chart(slide, x, y, w, h, *, series, x_title, y_title,
                          x_min=0, x_max=None, x_unit=None,
                          y_min=0, y_max=None, y_unit=None,
                          legend=False, legend_pos=None, smooth=False):
    """XY-scatter line+markers chart with the deck's standard styling.

    Use when you need data points to sit at arbitrary X-positions (e.g.,
    plotting MPL at the midpoint of each L-interval) while keeping tick
    marks at standard L-values.  Both axes are numeric value axes.

    series: list of ``(name, [(x, y), ...], color: RGBColor, marker: str)``.
    smooth: True for XY_SCATTER_SMOOTH (cubic spline through points),
            False for XY_SCATTER_LINES (straight segments).
    """
    _add_graphicframe_shadow(slide, x, y, w, h)
    cd = XyChartData()
    for name, points, _color, _marker in series:
        s = cd.add_series(name)
        for px, py in points:
            s.add_data_point(px, py)
    chart_type = (XL_CHART_TYPE.XY_SCATTER_SMOOTH if smooth
                  else XL_CHART_TYPE.XY_SCATTER_LINES)
    chart_shape = slide.shapes.add_chart(
        chart_type,
        int(x), int(y), int(w), int(h), cd,
    )
    chart = chart_shape.chart
    chart.has_title = False
    chart.has_legend = legend

    # Per-series styling
    for idx, ser in enumerate(chart.series):
        name, points, color, marker = series[idx]
        line = ser.format.line
        line.color.rgb = color
        line.width = Pt(2.5)
        ser_xml = ser._element
        clr_hex = f'{color[0]:02X}{color[1]:02X}{color[2]:02X}'
        for old in ser_xml.findall(qn('c:marker')):
            ser_xml.remove(old)
        m = ET.SubElement(ser_xml, qn('c:marker'))
        sym = ET.SubElement(m, qn('c:symbol')); sym.set('val', marker)
        sz_el = ET.SubElement(m, qn('c:size')); sz_el.set('val', '7')
        sp = ET.SubElement(m, qn('c:spPr'))
        fl = ET.SubElement(sp, qn('a:solidFill'))
        rg = ET.SubElement(fl, qn('a:srgbClr')); rg.set('val', clr_hex)
        ln = ET.SubElement(sp, qn('a:ln'))
        lf = ET.SubElement(ln, qn('a:solidFill'))
        lr = ET.SubElement(lf, qn('a:srgbClr')); lr.set('val', clr_hex)
        for sm in ser_xml.findall(qn('c:smooth')):
            ser_xml.remove(sm)
        smooth = ET.SubElement(ser_xml, qn('c:smooth'))
        smooth.set('val', '0')

    # Both axes are value axes in XY scatter; python-pptx returns the
    # X axis through .category_axis (wrapped as a ValueAxis since
    # catAx_lst is empty) and the Y axis through .value_axis.
    x_ax = chart.category_axis
    x_ax.tick_labels.font.name = "Calibri"
    x_ax.tick_labels.font.size = Pt(10)
    x_ax.tick_labels.font.color.rgb = NAVY
    if x_max is not None:
        x_ax.minimum_scale = x_min
        x_ax.maximum_scale = x_max
    if x_unit is not None:
        x_ax.major_unit = x_unit
    x_ax.has_title = True
    x_ax.axis_title.text_frame.text = x_title
    ar = x_ax.axis_title.text_frame.paragraphs[0].runs[0]
    ar.font.name = "Calibri"; ar.font.size = Pt(12)
    ar.font.bold = True; ar.font.italic = True; ar.font.color.rgb = NAVY

    y_ax = chart.value_axis
    y_ax.tick_labels.font.name = "Calibri"
    y_ax.tick_labels.font.size = Pt(10)
    y_ax.tick_labels.font.color.rgb = NAVY
    if y_max is not None:
        y_ax.minimum_scale = y_min
        y_ax.maximum_scale = y_max
    if y_unit is not None:
        y_ax.major_unit = y_unit
    y_ax.has_title = True
    y_ax.axis_title.text_frame.text = y_title
    ar = y_ax.axis_title.text_frame.paragraphs[0].runs[0]
    ar.font.name = "Calibri"; ar.font.size = Pt(12)
    ar.font.bold = True; ar.font.italic = True; ar.font.color.rgb = NAVY

    _add_dashed_gridlines(x_ax._element)
    _add_dashed_gridlines(y_ax._element)

    # Legend
    if legend and legend_pos is not None:
        leg_el = chart.legend._element
        chart.legend.font.name = "Calibri"
        chart.legend.font.size = Pt(11)
        chart.legend.font.color.rgb = NAVY
        chart.legend.include_in_layout = False
        for old in leg_el.findall(qn('c:layout')):
            leg_el.remove(old)
        for old in leg_el.findall(qn('c:legendPos')):
            leg_el.remove(old)
        pos = ET.SubElement(leg_el, qn('c:legendPos')); pos.set('val', 'tr')
        leg_el.remove(pos); leg_el.insert(0, pos)
        layout = ET.Element(qn('c:layout'))
        ml = ET.SubElement(layout, qn('c:manualLayout'))
        ET.SubElement(ml, qn('c:xMode')).set('val', 'edge')
        ET.SubElement(ml, qn('c:yMode')).set('val', 'edge')
        ET.SubElement(ml, qn('c:x')).set('val', legend_pos[0])
        ET.SubElement(ml, qn('c:y')).set('val', legend_pos[1])
        ET.SubElement(ml, qn('c:w')).set('val', legend_pos[2])
        ET.SubElement(ml, qn('c:h')).set('val', legend_pos[3])
        pos.addnext(layout)
        for old in leg_el.findall(qn('c:spPr')):
            leg_el.remove(old)
        leg_spPr = ET.Element(qn('c:spPr'))
        sf = ET.SubElement(leg_spPr, qn('a:solidFill'))
        clr = ET.SubElement(sf, qn('a:srgbClr')); clr.set('val', 'FFFFFF')
        ln = ET.SubElement(leg_spPr, qn('a:ln')); ln.set('w', '6350')
        lf = ET.SubElement(ln, qn('a:solidFill'))
        lc = ET.SubElement(lf, qn('a:srgbClr')); lc.set('val', '0B2B4E')
        layout.addnext(leg_spPr)

    return chart_shape


def _omml_acc_overline(symbol):
    """Inline OMML accent (overline / bar) on a math symbol.

    symbol: e.g. 'K' or 'L' – the variable to wear the bar.
    Returns an OMML fragment intended to be embedded inside an <m:oMath>.
    """
    return (
        '<m:acc>'
          '<m:accPr><m:chr m:val="̅"/></m:accPr>'
          '<m:e>'
            '<m:r>'
              '<a:rPr lang="en-US" b="0" i="1">'
                '<a:latin typeface="Cambria Math"/>'
              '</a:rPr>'
              f'<m:t>{symbol}</m:t>'
            '</m:r>'
          '</m:e>'
        '</m:acc>'
    )


def _add_mixed_textbox(slide, left, top, width, height, segments, *,
                        align=PP_ALIGN.LEFT, default_color=NAVY,
                        default_size=24,
                        margin_left=None, margin_right=None,
                        margin_top=None, margin_bottom=None):
    """Build a textbox whose paragraphs mix plain text runs and inline OMML.

    segments: list of (kind, content, opts) tuples, with kind ∈ {"text",
    "omml", "break"}.  "break" inserts a new paragraph.  Opts may set
    `size`, `bold`, `italic`, `color`, `font` per run.
    """
    box = slide.shapes.add_textbox(int(left), int(top),
                                     int(width), int(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05) if margin_left is None else margin_left
    tf.margin_right = Inches(0.05) if margin_right is None else margin_right
    tf.margin_top = Inches(0.0) if margin_top is None else margin_top
    tf.margin_bottom = Inches(0.0) if margin_bottom is None else margin_bottom

    align_attr = ''
    if align == PP_ALIGN.CENTER: align_attr = ' algn="ctr"'
    elif align == PP_ALIGN.RIGHT: align_attr = ' algn="r"'

    def _start_para():
        return [f'<a:p xmlns:a="{A_NS}" xmlns:m="{M_NS}" xmlns:a14="{A14_NS}">',
                f'<a:pPr{align_attr}/>' if align_attr else '']

    paragraphs = [_start_para()]
    for kind, content, opts in segments:
        if kind == 'break':
            paragraphs[-1].append('<a:endParaRPr lang="en-US"/></a:p>')
            paragraphs.append(_start_para())
            continue
        size_pt = int(opts.get('size', default_size) * 100)
        color = opts.get('color', default_color)
        clr_hex = f'{color[0]:02X}{color[1]:02X}{color[2]:02X}'
        bold_attr = ' b="1"' if opts.get('bold') else ''
        italic_attr = ' i="1"' if opts.get('italic') else ''
        font = opts.get('font', 'Calibri')
        if kind == 'text':
            paragraphs[-1].append(
                f'<a:r><a:rPr lang="en-US" sz="{size_pt}"{bold_attr}{italic_attr}>'
                f'<a:solidFill><a:srgbClr val="{clr_hex}"/></a:solidFill>'
                f'<a:latin typeface="{font}"/>'
                f'</a:rPr><a:t>{content}</a:t></a:r>'
            )
        elif kind == 'omml':
            paragraphs[-1].append(
                f'<a14:m><m:oMath>{content}</m:oMath></a14:m>'
            )
    paragraphs[-1].append('<a:endParaRPr lang="en-US"/></a:p>')
    full_xml = ''.join(''.join(p) for p in paragraphs)

    txBody = tf._txBody
    for old in list(txBody.findall(qn('a:p'))):
        txBody.remove(old)
    # We have to parse each <a:p> separately so the namespaces resolve
    for p_str in full_xml.split('</a:p>'):
        if not p_str.strip(): continue
        p_xml = p_str + '</a:p>'
        new_p = ET.fromstring(p_xml)
        txBody.append(new_p)
        # Apply size+color to any OMML m:r elements inside this paragraph.
        # Color is set to ``default_color`` ONLY when no per-run solidFill
        # is already present — this lets callers tint individual OMML runs
        # via the optional ``color=`` argument on ``_omml_run`` / ``_omml_text``
        # (e.g., the green ΔL / ΔQ in the slide-14 Convention box) without
        # being silently overridden here.
        clr_hex = f'{default_color[0]:02X}{default_color[1]:02X}{default_color[2]:02X}'
        for r in new_p.iter(qn('m:r')):
            arPr = r.find(qn('a:rPr'))
            if arPr is None:
                arPr = ET.Element(qn('a:rPr'))
                arPr.set('lang', 'en-US')
                r.insert(0, arPr)
            arPr.set('sz', str(int(default_size * 100)))
            if arPr.find(qn('a:solidFill')) is None:
                # fill BEFORE a:latin (schema order) or the color is ignored
                sf = arPr.makeelement(qn('a:solidFill'), {})
                srgb = ET.SubElement(sf, qn('a:srgbClr'))
                srgb.set('val', clr_hex)
                arPr.insert(0, sf)
    return box


def _add_math_equation(slide, left, top, width, height, omml_content, *,
                       size_pt=32, color=NAVY, fill=None, line=None,
                       rounded=False, shadow=False, corner_pct=25000):
    """Place an OMML equation in a textbox on the slide.

    omml_content: a string built from _omml_* helpers (without the outer
    <m:oMathPara> wrapper).

    ``rounded=True`` gives the fill box rounded corners (``corner_pct`` in
    OOXML adj units, 25000 ≈ 25%); ``shadow=True`` adds a soft drop
    shadow.  Both default off.  IMPORTANT: the roundRect prstGeom must be
    inserted right after <a:xfrm> and BEFORE the fill, or PowerPoint
    silently drops the shape (see slide-54 fix).
    """
    left, top, width, height = int(left), int(top), int(width), int(height)
    box = slide.shapes.add_textbox(left, top, width, height)

    if fill is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
    if line is not None:
        box.line.color.rgb = line
        box.line.width = Pt(0.75)

    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Replace the default <a:p> with one that hosts a math zone (<a14:m>)
    # containing the oMathPara.  The <a14:m> wrapper is REQUIRED by
    # PowerPoint to recognise OMML inside a textbox – without it PPT just
    # shows empty boxes.
    txBody = tf._txBody
    for p in list(txBody.findall(qn('a:p'))):
        txBody.remove(p)

    sz = int(size_pt * 100)
    clr_hex = '{:02X}{:02X}{:02X}'.format(color[0], color[1], color[2])

    # Use unique namespace prefixes; lxml will resolve them when parsing.
    p_xml = (
        f'<a:p xmlns:a="{A_NS}" xmlns:m="{M_NS}" xmlns:a14="{A14_NS}">'
        f'<a:pPr algn="ctr">'
        f'<a:defRPr sz="{sz}" b="0" i="1">'
        f'<a:solidFill><a:srgbClr val="{clr_hex}"/></a:solidFill>'
        f'<a:latin typeface="Cambria Math"/>'
        f'</a:defRPr>'
        f'</a:pPr>'
        f'<a14:m>'
        f'<m:oMathPara>'
        f'<m:oMathParaPr><m:jc m:val="centerGroup"/></m:oMathParaPr>'
        f'<m:oMath>{omml_content}</m:oMath>'
        f'</m:oMathPara>'
        f'</a14:m>'
        f'<a:endParaRPr lang="en-US" sz="{sz}"/>'
        f'</a:p>'
    )

    new_p = ET.fromstring(p_xml)
    txBody.append(new_p)

    # Set size and color on every OMML run's <a:rPr> so the text renders at
    # the right size.  (The defRPr above is the fallback.)
    for r in new_p.iter(qn('m:r')):
        arPr = r.find(qn('a:rPr'))
        if arPr is None:
            arPr = ET.Element(qn('a:rPr'))
            r.insert(0, arPr)
        arPr.set('sz', str(sz))
        if arPr.get('lang') is None:
            arPr.set('lang', 'en-US')
        # respect per-run colors set via _omml_run/_omml_text(color=...);
        # only runs WITHOUT a fill get the equation's default color
        # (was clobbering firm-color runs — Nico 2026-08-07)
        if not arPr.findall(qn('a:solidFill')):
            sf = ET.Element(qn('a:solidFill'))
            srgb = ET.SubElement(sf, qn('a:srgbClr'))
            srgb.set('val', clr_hex)
            arPr.insert(0, sf)

    # Optional rounded corners + drop shadow on the fill box.  prstGeom
    # MUST sit right after <a:xfrm> and before <a:solidFill>, else PPT
    # silently refuses to render the shape (slide-54 lesson).
    if rounded or shadow:
        spPr = box._element.find(qn('p:spPr'))
        if spPr is not None:
            if rounded:
                for old in spPr.findall(qn('a:prstGeom')):
                    spPr.remove(old)
                prstGeom = ET.Element(qn('a:prstGeom'))
                prstGeom.set('prst', 'roundRect')
                avLst = ET.SubElement(prstGeom, qn('a:avLst'))
                gd = ET.SubElement(avLst, qn('a:gd'))
                gd.set('name', 'adj'); gd.set('fmla', f'val {int(corner_pct)}')
                xfrm = spPr.find(qn('a:xfrm'))
                if xfrm is not None:
                    xfrm.addnext(prstGeom)
                else:
                    spPr.insert(0, prstGeom)
            if shadow:
                _add_drop_shadow(box)
    return box


CREAM = RGBColor(0xFD, 0xF6, 0xE6)


DIM = RGBColor(0xBF, 0xBF, 0xBF)           # 2026-08-24 (Nico): outline


                                           # items not currently covered
                                           # are shaded (his video decks
                                           # use schemeClr bg1 lumMod 75%
                                           # over white = #BFBFBF)
RED = RGBColor(0xC0, 0x00, 0x00)           # source red (C00000)


RED_FF = RGBColor(0xFF, 0x00, 0x00)        # source bright red (FF0000)


GREEN_DK = RGBColor(0x00, 0x7A, 0x33)
EMC_PURPLE = RGBColor(0x70, 0x30, 0xA0)   # Nico's external-marginal-cost purple (2026-08-30)      # the deck's shift green


                                           # (2026-08-23, Nico: the
                                           # source bright green was
                                           # too light everywhere)
BLUE_PED = RGBColor(0x00, 0x70, 0xC0)      # source blue (S', excess supply)


STEEL = RGBColor(0x95, 0xB3, 0xD7)         # source light-blue supply curve


DARKRED = RGBColor(0xA2, 0x16, 0x2A)       # source MC bar red


NB_BLUE = RGBColor(0x2E, 0x5B, 0x9F)       # net-benefit bar blue


def _add_slidenum_field(slide, left, top, width, height, page_num, *,
                        size=12, color=GRAY):
    tb = slide.shapes.add_textbox(int(left), int(top), int(width), int(height))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    p_el = p._p
    guid = str(uuid.uuid5(uuid.NAMESPACE_DNS, f"m1rev-slidenum-{page_num}")).upper()
    fld = p_el.makeelement(qn('a:fld'), {'id': '{%s}' % guid, 'type': 'slidenum'})
    rPr = p_el.makeelement(qn('a:rPr'),
                           {'lang': 'en-US', 'sz': str(size * 100), 'dirty': '0'})
    fill = p_el.makeelement(qn('a:solidFill'), {})
    clr = p_el.makeelement(qn('a:srgbClr'), {'val': str(color)})
    fill.append(clr)
    rPr.append(fill)
    latin = p_el.makeelement(qn('a:latin'), {'typeface': 'Calibri'})
    rPr.append(latin)
    t = p_el.makeelement(qn('a:t'), {})
    t.text = str(page_num)
    fld.append(rPr)
    fld.append(t)
    p_el.append(fld)
    return tb


def _draw_footer(slide, footer_text, page_num):  # noqa: F811 — field override
    _add_rect(slide, 0, Inches(7.15), SLIDE_W, Inches(0.02), RULE)
    _add_rect(slide, MARGIN, Inches(7.135), GOLD_W, Inches(0.05), GOLD)
    _add_text(slide, MARGIN, Inches(7.20), Inches(11), Inches(0.32),
              footer_text, size=12, color=GRAY)
    _add_slidenum_field(slide, Inches(12.55), Inches(7.20), Inches(0.55),
                        Inches(0.32), page_num)


def _add_media_image(slide, fname, *, left, top, width=None, height=None,
                     rounded=True, shadow=True, corner_pct=8):
    """Place a source-deck image by filename from _source_images/.
    Logos / screenshots / clippings: rounded=False, shadow=False."""
    path = SRC_IMG_DIR / fname
    kwargs = {"left": int(left), "top": int(top)}
    if width is not None:
        kwargs["width"] = int(width)
    if height is not None:
        kwargs["height"] = int(height)
    pic = slide.shapes.add_picture(str(path), **kwargs)
    if rounded:
        _apply_picture_style(pic, corner_pct=corner_pct)
    elif shadow:
        _add_drop_shadow(pic)
    return pic


def _set_cell_borders(cell, *, color=RULE, weight_pt=0.75):
    tcPr = cell._tc.get_or_add_tcPr()
    for tag in ('a:lnL', 'a:lnR', 'a:lnT', 'a:lnB'):
        for old in tcPr.findall(qn(tag)):
            tcPr.remove(old)
    for tag in reversed(('a:lnL', 'a:lnR', 'a:lnT', 'a:lnB')):
        ln = tcPr.makeelement(qn(tag), {'w': str(int(weight_pt * 12700)),
                                        'cap': 'flat'})
        fill = ln.makeelement(qn('a:solidFill'), {})
        c = fill.makeelement(qn('a:srgbClr'), {'val': str(color)})
        fill.append(c)
        ln.append(fill)
        tcPr.insert(0, ln)


def _add_styled_table(slide, left, top, width, height, rows_data, *,
                      col_widths=None, row_heights=None, font_size=18,
                      header_size=None, backing_pad=Inches(0.15),
                      first_col_bold=True, first_col_align_left=True,
                      cell_fills=None, cell_text_colors=None):
    """rows_data: list of rows (row 0 = navy header). cell_fills /
    cell_text_colors: optional {(r, c): RGBColor} overrides."""
    header_size = header_size or font_size
    left, top, width, height = int(left), int(top), int(width), int(height)
    _add_graphicframe_shadow(slide, left - int(backing_pad),
                             top - int(backing_pad),
                             width + 2 * int(backing_pad),
                             height + 2 * int(backing_pad))
    n_rows, n_cols = len(rows_data), len(rows_data[0])
    gf = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = gf.table
    tblPr = tbl._tbl.find(qn('a:tblPr'))
    if tblPr is not None:
        tblPr.set('firstRow', '0')
        tblPr.set('bandRow', '0')
        for child in list(tblPr):
            tblPr.remove(child)
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = int(w)
    if row_heights:
        for i, h in enumerate(row_heights):
            tbl.rows[i].height = int(h)
    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.margin_left = Inches(0.10)
            cell.margin_right = Inches(0.10)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            hdr = (r == 0)
            if cell_fills and (r, c) in cell_fills:
                cell.fill.solid()
                cell.fill.fore_color.rgb = cell_fills[(r, c)]
            elif hdr:
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else CREAM
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            if first_col_align_left and c == 0 and not hdr:
                p.alignment = PP_ALIGN.LEFT
            else:
                p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(val)
            run.font.name = "Calibri"
            run.font.size = Pt(header_size if hdr else font_size)
            run.font.bold = hdr or (first_col_bold and c == 0)
            if cell_text_colors and (r, c) in cell_text_colors:
                run.font.color.rgb = cell_text_colors[(r, c)]
            else:
                run.font.color.rgb = WHITE if hdr else NAVY
            _set_cell_borders(cell)
    return gf


# Poll Break badge — the ONE geometry from Teaching CLAUDE.md
# ("The Poll Break parallelogram has ONE fixed geometry"), 2026-08-27.
# Built from these numbers rather than copied from a hand-tuned sidecar,
# so the rule alone is enough to regenerate it.
POLLBREAK_XY = (9361444, 6190030)        # EMU: left 10.2377", top 6.7695"


POLLBREAK_WH = (2697480, 487009)         # EMU: 2.9500 x 0.5326"


POLLBREAK_SLANT = 0.72 / 2.95            # 0.72" slant on a 2.95" width


# custGeom path in the normalised 100000 x 100000 box: s = the slant as a
# fraction of the width, r = the corner rounding, d = r*s/100000.
_PB_S = int(round(POLLBREAK_SLANT * 100000))     # 24406


_PB_R = 5000


_PB_D = int(round(_PB_R * _PB_S / 100000.0))     # 1220


_PB_PATH = (
    '<a:custGeom><a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/>'
    '<a:rect l="{s}" t="0" r="{ws}" b="100000"/>'
    '<a:pathLst><a:path w="100000" h="100000">'
    '<a:moveTo><a:pt x="{sr}" y="0"/></a:moveTo>'
    '<a:lnTo><a:pt x="{wr}" y="0"/></a:lnTo>'
    '<a:cubicBezTo><a:pt x="100000" y="0"/><a:pt x="100000" y="0"/>'
    '<a:pt x="{wd}" y="{r}"/></a:cubicBezTo>'
    '<a:lnTo><a:pt x="{wsd}" y="{hr}"/></a:lnTo>'
    '<a:cubicBezTo><a:pt x="{ws}" y="100000"/><a:pt x="{ws}" y="100000"/>'
    '<a:pt x="{wsr}" y="100000"/></a:cubicBezTo>'
    '<a:lnTo><a:pt x="{r}" y="100000"/></a:lnTo>'
    '<a:cubicBezTo><a:pt x="0" y="100000"/><a:pt x="0" y="100000"/>'
    '<a:pt x="{d}" y="{hr}"/></a:cubicBezTo>'
    '<a:lnTo><a:pt x="{sd}" y="{r}"/></a:lnTo>'
    '<a:cubicBezTo><a:pt x="{s}" y="0"/><a:pt x="{s}" y="0"/>'
    '<a:pt x="{sr}" y="0"/></a:cubicBezTo>'
    '<a:close/></a:path></a:pathLst></a:custGeom>'
).format(s=_PB_S, r=_PB_R, d=_PB_D, sr=_PB_S + _PB_R, sd=_PB_S - _PB_D,
         wr=100000 - _PB_R, wd=100000 - _PB_D, ws=100000 - _PB_S,
         wsr=100000 - _PB_S - _PB_R, wsd=100000 - _PB_S + _PB_D,
         hr=100000 - _PB_R)


_PB_SHADOW = ('<a:effectLst><a:outerShdw blurRad="50800" dist="38100" '
              'dir="2700000" algn="tl" rotWithShape="0">'
              '<a:srgbClr val="000000"><a:alpha val="50000"/></a:srgbClr>'
              '</a:outerShdw></a:effectLst>')


_PB_XML = (
    '<p:grpSp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
    'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
    '<p:nvGrpSpPr><p:cNvPr id="9600" name="PollBreakBadge"/>'
    '<p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>'
    '<p:grpSpPr><a:xfrm><a:off x="{x}" y="{y}"/><a:ext cx="{cx}" cy="{cy}"/>'
    '<a:chOff x="{x}" y="{y}"/><a:chExt cx="{cx}" cy="{cy}"/></a:xfrm></p:grpSpPr>'
    # the gold parallelogram
    '<p:sp><p:nvSpPr><p:cNvPr id="9601" name="PollBreakShape"/>'
    '<p:cNvSpPr/><p:nvPr/></p:nvSpPr><p:spPr>'
    '<a:xfrm><a:off x="{x}" y="{y}"/><a:ext cx="{cx}" cy="{cy}"/></a:xfrm>'
    '{path}<a:solidFill><a:srgbClr val="E09F3E"/></a:solidFill>'
    '<a:ln><a:noFill/></a:ln>{shadow}</p:spPr>'
    '<p:txBody><a:bodyPr rtlCol="0" anchor="ctr"/><a:lstStyle/>'
    '<a:p><a:pPr algn="ctr"/><a:endParaRPr/></a:p></p:txBody></p:sp>'
    # the navy label, filling the parallel middle (W - 2S)
    '<p:sp><p:nvSpPr><p:cNvPr id="9602" name="PollBreakLabel"/>'
    '<p:cNvSpPr txBox="1"/><p:nvPr/></p:nvSpPr><p:spPr>'
    '<a:xfrm><a:off x="{lx}" y="{y}"/><a:ext cx="{lcx}" cy="{cy}"/></a:xfrm>'
    '<a:prstGeom prst="rect"><a:avLst/></a:prstGeom><a:noFill/></p:spPr>'
    '<p:txBody><a:bodyPr wrap="none" lIns="45720" tIns="18288" rIns="45720" '
    'bIns="18288" anchor="ctr"><a:spAutoFit/></a:bodyPr><a:lstStyle/>'
    '<a:p><a:pPr algn="ctr"/><a:r><a:rPr sz="2800" b="1">'
    '<a:solidFill><a:srgbClr val="0B2B4E"/></a:solidFill>'
    '<a:latin typeface="Calibri"/></a:rPr><a:t>Poll Break</a:t></a:r></a:p>'
    '</p:txBody></p:sp></p:grpSp>'
).format(x=POLLBREAK_XY[0], y=POLLBREAK_XY[1],
         cx=POLLBREAK_WH[0], cy=POLLBREAK_WH[1],
         lx=POLLBREAK_XY[0] + int(round(POLLBREAK_WH[0] * POLLBREAK_SLANT)),
         lcx=POLLBREAK_WH[0] - 2 * int(round(POLLBREAK_WH[0] * POLLBREAK_SLANT)),
         path=_PB_PATH, shadow=_PB_SHADOW)


def _add_pollbreak_badge(slide, _ids=[9600]):
    """The deck-standard Poll Break badge, bottom-right, IN FRONT of the
    footer (it straddles the footer rule at y 7.15", so it must be the
    LAST shape appended). Call AFTER _draw_footer."""
    el = ET.fromstring(_PB_XML)
    for i, nv in enumerate(el.iter(qn('p:cNvPr'))):
        nv.set('id', str(_ids[0] + i))
    _ids[0] += 10
    slide.shapes._spTree.append(el)
    return el


def make_stub(prs, page_num, section_tag, title, note, *, hidden=False):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, section_tag)
    _draw_action_title(slide, title)
    _add_text(slide, MARGIN, Inches(3.4), RULE_W, Inches(0.8),
              "[ %s ]" % note, size=20, italic=True, color=GRAY,
              font="Calibri", align=PP_ALIGN.CENTER)
    _draw_footer(slide, FOOTER_TEXT, page_num)
    if hidden:
        slide._element.set('show', '0')
    return slide


class SimpleFig:
    def __init__(self, left_in, bottom_in, w_in, h_in, xmax, ymax):
        self.l, self.b, self.w, self.h = left_in, bottom_in, w_in, h_in
        self.xmax, self.ymax = xmax, ymax

    def x(self, xv):
        return Inches(self.l + self.w * xv / self.xmax)

    def y(self, yv):
        return Inches(self.b - self.h * yv / self.ymax)


# Clearance between the y-axis line and the right edge of its title box
# (inches).  Hand-set by Nico on slide 36, 2026-08-30.
Y_TITLE_GAP = 0.08


def _text_w_in(text, size_pt, *, bold=False, italic=False):
    """Rendered width of `text` in INCHES, measured in the real Calibri face.

    Axis titles are sized tight to their label (Teaching CLAUDE.md), so the
    anchor arithmetic needs the true width rather than a guess.  Falls back
    to a rough estimate if the font file is not present.
    """
    face = {(False, False): "calibri.ttf", (True, False): "calibrib.ttf",
            (False, True): "calibrii.ttf", (True, True): "calibriz.ttf"}[
                (bool(bold), bool(italic))]
    try:
        from PIL import ImageFont
        f = ImageFont.truetype("C:/Windows/Fonts/" + face,
                               int(round(size_pt * 96 / 72.0)))
        b = f.getbbox(text)
        return (b[2] - b[0]) / 96.0
    except Exception:
        return 0.62 * len(text) * size_pt / 72.0


def _fig_axes(slide, fig, *, weight_pt=2.0,
              x_title="Quantity", y_title="Price ($)", label_size=18,
              titles_at_tip=True):
    """Navy arrow axes + axis titles (y above the top arrow, x right of
    the right arrow, per the source-chart convention)."""
    _add_arrow(slide, (fig.x(0), fig.y(0)),
               (fig.x(0), Inches(fig.b - fig.h - 0.18)),
               color=NAVY, weight_pt=weight_pt, head=True)
    _add_arrow(slide, (fig.x(0), fig.y(0)),
               (Inches(fig.l + fig.w + 0.18), fig.y(0)),
               color=NAVY, weight_pt=weight_pt, head=True)
    # titles_at_tip: the y title sits just LEFT of the y-arrow's tip and
    # level with it, the x title just BELOW the x-arrow's tip (Teaching
    # CLAUDE.md, 2026-08-30).  This is the DECK-WIDE default as of
    # 2026-08-30; pass False only for a chart that needs the old placement.
    if y_title:
        if titles_at_tip:
            # Box sized TIGHT to the label; its vertical MIDDLE sits exactly
            # at the arrow tip, and its right border Y_TITLE_GAP clear of the
            # axis line -- sitting ON the axis read as too tight.  (Gap
            # measured off Nico's recalibration on slide 36, 2026-08-30:
            # y-axis at x 3.050, "P" box right edge at 2.970.)
            w = _text_w_in(y_title, label_size, bold=True,
                           italic=True) + 0.08
            tip_y = fig.b - fig.h - 0.18
            _add_text(slide, Inches(fig.l - w - Y_TITLE_GAP),
                      Inches(tip_y - 0.145),
                      Inches(w), Inches(0.29), y_title, size=label_size,
                      bold=True, italic=True, color=NAVY, font="Calibri")
        else:
            _add_text(slide, Inches(fig.l - 0.75),
                      Inches(fig.b - fig.h - 0.62),
                      Inches(2.0), Inches(0.32), y_title, size=label_size,
                      bold=True, italic=True, color=NAVY, font="Calibri")
    if x_title:
        if titles_at_tip:
            # Box sized tight; its horizontal MIDPOINT sits exactly at the
            # x-arrow's tip, and its top 0.05" under the axis.
            w = _text_w_in(x_title, label_size, bold=True,
                           italic=True) + 0.08
            tip_x = fig.l + fig.w + 0.18
            _add_text(slide, Inches(tip_x - w / 2.0), Inches(fig.b + 0.05),
                      Inches(w), Inches(0.29), x_title, size=label_size,
                      bold=True, italic=True, color=NAVY, font="Calibri")
        else:
            _add_text(slide, Inches(fig.l + fig.w - 0.4),
                      Inches(fig.b + 0.10),
                      Inches(1.8), Inches(0.32), x_title, size=label_size,
                      bold=True, italic=True, color=NAVY, font="Calibri")


def _fig_line(slide, fig, p0, p1, *, color=NAVY, weight_pt=2.5, dash=None,
              head=False):
    """Straight curve segment in logical coords."""
    return _add_arrow(slide, (fig.x(p0[0]), fig.y(p0[1])),
                      (fig.x(p1[0]), fig.y(p1[1])),
                      color=color, weight_pt=weight_pt, dash=dash, head=head)


def _fig_guide(slide, fig, pt, *, color=GRAY, to_x=True, to_y=True,
               weight_pt=1.25, dash='dash'):
    """Dashed guide lines from a point to both axes. Returns
    (horizontal, vertical) so callers can name/group them."""
    x, y = pt
    h = v = None
    if to_y:
        h = _add_arrow(slide, (fig.x(0), fig.y(y)), (fig.x(x), fig.y(y)),
                       color=color, weight_pt=weight_pt, head=False,
                       dash=dash)
    if to_x:
        v = _add_arrow(slide, (fig.x(x), fig.y(y)), (fig.x(x), fig.y(0)),
                       color=color, weight_pt=weight_pt, head=False,
                       dash=dash)
    return h, v


def _fig_ylab(slide, fig, yv, label, *, color=NAVY, size=18, bold=False,
              width_in=None):
    """A y-axis tick label, right-aligned against the axis.

    2026-08-30 (Nico): the box is sized to the label, so it re-fits itself
    when the text changes.  The right edge - and so the rendered text - is
    unchanged.  Pass width_in only to pin a box."""
    if width_in is None:
        width_in = _text_w_in(label, size, bold=bold, italic=True) + 0.08
    return _add_text(slide, Inches(fig.l - width_in - 0.08),
                     fig.y(yv) - Inches(0.15),
                     Inches(width_in), Inches(0.3), label, size=size,
                     bold=bold, italic=True, color=color, font="Calibri",
                     align=PP_ALIGN.RIGHT)


def _fig_ylab_subsup(slide, fig, yv, base, sub, sup, *, color=NAVY,
                     size=16, bold=True):
    """A y-axis tick label carrying a subscript AND a superscript, e.g.
    P with a subscript E and a superscript C (2026-08-30, Nico — the
    price paid by consumers once the externality is taxed).  Right-aligned
    against the axis, in a box sized to the three runs."""
    w = (_text_w_in(base, size, bold=bold, italic=True)
         + _text_w_in(sub, size * 0.72, bold=bold, italic=True)
         + _text_w_in(sup, size * 0.72, bold=bold, italic=True) + 0.10)
    box = _add_text(slide, Inches(fig.l - w - 0.08),
                    fig.y(yv) - Inches(0.15), Inches(w), Inches(0.30), "",
                    size=size, color=color, font="Calibri",
                    align=PP_ALIGN.RIGHT)
    p = box.text_frame.paragraphs[0]
    for text, base_ln in ((base, None), (sub, "-25000"), (sup, "30000")):
        r = p.add_run()
        r.text = text
        r.font.name = "Calibri"
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = True
        r.font.color.rgb = color
        if base_ln:
            r._r.find(qn("a:rPr")).set("baseline", base_ln)
    return box


def _fig_xlab(slide, fig, xv, label, *, color=NAVY, size=18, bold=False):
    """An x-axis tick label, centred on the tick in a box sized to it."""
    w = _text_w_in(label, size, bold=bold, italic=True) + 0.08
    return _add_text(slide, fig.x(xv) - Inches(w / 2), Inches(fig.b + 0.06),
                     Inches(w), Inches(0.3), label, size=size,
                     bold=bold, italic=True, color=color, font="Calibri",
                     align=PP_ALIGN.CENTER)


def _fig_vbrace(slide, fig, y_lo, y_hi, x_dev, label, *, color=NAVY,
                size=14, depth=0.16, gap=0.08, label_w=1.55):
    """A vertical brace spanning y_lo..y_hi at x_dev inches, opening to the
    RIGHT (its point faces right, toward the figure), with `label` to its
    left.  Used for Nico's P_B - P_S = t bracket on slide 64."""
    y0, y1 = fig.y(y_hi), fig.y(y_lo)
    span = y1 - y0
    shp = slide.shapes.add_shape(
        MSO_SHAPE.LEFT_BRACE, int(Inches(x_dev)), int(y0),
        int(Inches(depth)), int(span))
    shp.fill.background()
    shp.line.color.rgb = color
    shp.line.width = Pt(1.75)
    shp.shadow.inherit = False
    # the caller may place the label itself (pass label=None) when the
    # left margin is too narrow to hold it
    if label:
        _add_text(slide, int(Inches(x_dev - gap) - Inches(label_w)),
                  int(y0 + span / 2 - Inches(0.17)),
                  int(Inches(label_w)), Inches(0.34), label, size=size,
                  bold=True, color=color, font="Calibri",
                  align=PP_ALIGN.RIGHT)
    return shp


def _fig_underbrace(slide, fig, x_lo, x_hi, y_dev, label, *,
                    color=NAVY, size=14, depth=0.18, gap=0.06):
    """A horizontal under-brace spanning x_lo..x_hi with `label` beneath.

    Built from a LEFT_BRACE rotated 270 degrees, so its point faces DOWN;
    PowerPoint rotates about the shape centre, so the un-rotated shape is
    authored `depth` wide by `span` tall and centred on where the brace
    should sit.  `y_dev` is the brace's vertical centre, in inches.
    """
    x0, x1 = fig.x(x_lo), fig.x(x_hi)
    span = x1 - x0
    cx, cy = (x0 + x1) / 2.0, int(Inches(y_dev))
    shp = slide.shapes.add_shape(
        MSO_SHAPE.LEFT_BRACE, int(cx - Inches(depth) / 2),
        int(cy - span / 2), int(Inches(depth)), int(span))
    shp.rotation = 270
    shp.fill.background()
    shp.line.color.rgb = color
    shp.line.width = Pt(1.75)
    shp.shadow.inherit = False
    # the label needs its own width - a short span would wrap it
    lw = max(span, int(Inches(3.0)))
    _add_text(slide, int(cx - lw / 2),
              int(cy + Inches(depth) / 2 + Inches(gap)),
              int(lw), Inches(0.34), label, size=size, bold=True,
              color=color, font="Calibri", align=PP_ALIGN.CENTER)
    return shp


_SWATCH_GEOM = {"sq": (MSO_SHAPE.RECTANGLE, False),
                "tri": (MSO_SHAPE.RIGHT_TRIANGLE, False),
                "tri_v": (MSO_SHAPE.RIGHT_TRIANGLE, True)}


def _welfare_rows(slide, left, top, width, rows, *, pitch=0.46,
                  size=19, swatch=0.20, alpha=30, text_dx=0.62):
    """The welfare-effects legend: each line preceded by a mark in the
    colour AND SHAPE of the region it names, and - where the line names a
    COMBINATION of regions - with those marks arranged the way the regions
    sit in the graph (2026-08-30, Nico).

    A row is ``(layout, marks, text)``:
      layout  "h" the marks sit side by side (areas side by side in the
              graph), "v" they stack (one area above the other);
      marks   a list of ``(colour, kind)``, kind in "sq" / "tri" / "tri_v"
              ("tri_v" is flipped, mirroring a region whose right angle is
              on top);
      text    the line itself, vertically centred on the mark block.
    """
    y = top
    for row in rows:
        layout, marks, text = row
        n = len(marks)
        step = swatch + 0.02
        block_h = swatch if layout == "h" else (n - 1) * (swatch + 0.01) + swatch
        for j, mark in enumerate(marks):
            col, kind = mark if isinstance(mark, tuple) else (mark, "sq")
            geom, flip = _SWATCH_GEOM[kind]
            dx = j * step if layout == "h" else 0.0
            dy = 0.0 if layout == "h" else j * (swatch + 0.01)
            sq = slide.shapes.add_shape(
                geom, int(left + Inches(dx)), int(Inches(y + dy)),
                int(Inches(swatch)), int(Inches(swatch)))
            if flip:
                xfrm = sq._element.find(".//" + qn("a:xfrm"))
                if xfrm is not None:
                    xfrm.set("flipV", "1")
            sq.fill.solid()
            sq.fill.fore_color.rgb = col
            _set_fill_alpha(sq, alpha)
            sq.line.color.rgb = NAVY
            sq.line.width = Pt(0.75)
            sq.shadow.inherit = False
        _add_text(slide, int(left + Inches(text_dx)),
                  int(Inches(y + block_h / 2.0 - 0.17)),
                  int(width - Inches(text_dx)),
                  Inches(0.34), text, size=size, color=NAVY,
                  font="Calibri")
        # a stacked row is taller than one line, so the next row clears it
        y += max(pitch, block_h + 0.06)


def _fig_curve_label(slide, fig, xv, yv, label, *, color=NAVY, size=20,
                     bold=True, width=None):
    """A curve label, in a box just wide enough for the text.

    2026-08-30 (Nico): every label in a figure - not only the axis titles
    - is sized to its own text, so it re-fits itself when the wording
    changes and never leaves a wide invisible box to catch the cursor.
    Pass `width` only to pin a label whose box must stay a fixed size.
    """
    w = width if width is not None else (
        _text_w_in(label, size, bold=bold, italic=True) + 0.08)
    return _add_text(slide, fig.x(xv), fig.y(yv) - Inches(0.16),
                     Inches(w), Inches(0.32), label, size=size, bold=bold,
                     italic=True, color=color, font="Calibri")


def content_slide(prs, page_num, section_tag, title, bullets, *,
                  size=24, sub_size=None, line_spacing_pts=12,
                  bullets_left=None, bullets_width=None,
                  bullets_top=Inches(1.6), bullets_height=Inches(5.35),
                  title_size=None, extras=None, notes=None):
    """M1 standard content slide. Bullet block is MIDDLE-anchored in the
    content region so the text sits vertically centered (course rule)."""
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, section_tag)
    if title_size is None:
        _draw_action_title(slide, title)
    else:
        _add_text(slide, MARGIN, Inches(0.55), RULE_W, Inches(0.7),
                  title, size=title_size, bold=True, color=NAVY,
                  font="Calibri")
        _add_rect(slide, MARGIN, Inches(1.25), RULE_W, Inches(0.02), RULE)
        _add_rect(slide, MARGIN, Inches(1.235), GOLD_W, Inches(0.05), GOLD)
    normalized = [(b, 0) if isinstance(b, str) else b for b in bullets]
    box = _add_hierarchical_bullets(
        slide,
        left=MARGIN if bullets_left is None else bullets_left,
        top=bullets_top,
        width=RULE_W if bullets_width is None else bullets_width,
        height=bullets_height,
        items=normalized,
        size=size, sub_size=sub_size,
        line_spacing_pts=line_spacing_pts,
    )
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    if extras is not None:
        extras(slide)
    _draw_footer(slide, FOOTER_TEXT, page_num)
    if notes:
        _set_notes(slide, notes)
    return slide


def _highlight_texts(slide, texts, color='FFFF00'):
    """Add <a:highlight> to every run whose text is in *texts*.
    Schema order inside rPr: fill < highlight < latin."""
    remaining = set(texts)
    for shp in slide.shapes:
        if not shp.has_text_frame:
            continue
        for para in shp.text_frame.paragraphs:
            for run in para.runs:
                if run.text in remaining:
                    rPr = run._r.get_or_add_rPr()
                    hl = rPr.makeelement(qn('a:highlight'), {})
                    srgb = hl.makeelement(qn('a:srgbClr'), {'val': color})
                    hl.append(srgb)
                    latin = rPr.find(qn('a:latin'))
                    if latin is not None:
                        latin.addprevious(hl)
                    else:
                        rPr.append(hl)


def _link_runs(slide, mapping):
    """Set external hyperlinks on runs by exact run text.
    mapping: {run_text: url}"""
    for shp in slide.shapes:
        if not shp.has_text_frame:
            continue
        for para in shp.text_frame.paragraphs:
            for run in para.runs:
                if run.text in mapping:
                    run.hyperlink.address = mapping[run.text]


# The preset shades its glyph from the shape FILL (darkenLess) and only
# strokes it with the line colour, so a WHITE face renders the triangle
# grey, not navy — verified 2026-08-23 against a render. An all-navy
# symbol therefore means a navy face with the glyph knocked out in white.
JUMP_BTN_FILL = NAVY          # backup-jump symbol: navy only (Nico)


JUMP_BTN_GLYPH = WHITE        # glyph knocked out of the navy face


EXT_BTN_FILL = NAVY           # 2026-08-23: external markers go navy too,


EXT_BTN_GLYPH = WHITE         # matching the jump buttons (Nico)


EXT_LINK_SHAPES = {
    "sound": MSO_SHAPE.ACTION_BUTTON_SOUND,        # podcast / audio
    "document": MSO_SHAPE.ACTION_BUTTON_DOCUMENT,  # article / paper
    "movie": MSO_SHAPE.ACTION_BUTTON_MOVIE,        # video
}


def _add_jump_button(slide, target_slide, *, left, top,
                     width=Inches(0.434), height=Inches(0.210), back=False):
    """Standalone action button carrying a jump-to-slide action.
    back=False -> points right (into the backup section);
    back=True  -> points left (unused since the back pills reverted to
    plain '← Back' text)."""
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ACTION_BUTTON_BEGINNING if back
        else MSO_SHAPE.ACTION_BUTTON_END,
        int(left), int(top), int(width), int(height))
    shp.fill.solid()
    shp.fill.fore_color.rgb = JUMP_BTN_FILL
    shp.line.color.rgb = JUMP_BTN_GLYPH
    shp.line.width = Pt(1.0)
    shp.shadow.inherit = False
    _add_drop_shadow(shp)
    shp.click_action.target_slide = target_slide
    return shp


def _add_ext_link_button(slide, kind, *, left, top, url=None,
                         width=Inches(0.60), height=Inches(0.30)):
    """External-link marker in the same action-button family as the
    slide-jump buttons, but keyed to WHAT it opens rather than to a
    direction — so an audio, an article and a video are told apart at a
    glance (2026-08-23, replaces the gold ▶ glyphs)."""
    shp = slide.shapes.add_shape(EXT_LINK_SHAPES[kind],
                                 int(left), int(top),
                                 int(width), int(height))
    shp.fill.solid()
    shp.fill.fore_color.rgb = EXT_BTN_FILL
    shp.line.color.rgb = EXT_BTN_GLYPH
    shp.line.width = Pt(1.25)
    shp.shadow.inherit = False
    _add_drop_shadow(shp)
    if url:
        shp.click_action.hyperlink.address = url
    return shp


def _add_jump_pill(slide, target_slide, *, left, top, width, label,
                   height=Inches(0.5), back=False, fill=None,
                   text_color=None, border=None, size=15):
    """Labelled jump affordance: a rounded pill carrying the label and the
    jump, with the action button seated in a reserved left inset (also
    clickable). Forward pills are white with a gold border and navy text;
    back pills are navy with white text."""
    fill = WHITE if fill is None else fill
    text_color = NAVY if text_color is None else text_color
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  int(left), int(top),
                                  int(width), int(height))
    try:
        pill.adjustments[0] = 0.28
    except Exception:
        pass
    pill.fill.solid()
    pill.fill.fore_color.rgb = fill
    if border is None:
        pill.line.fill.background()
    else:
        pill.line.color.rgb = border
        pill.line.width = Pt(1.5)
    pill.shadow.inherit = False
    _add_drop_shadow(pill)
    btn_w, btn_h = Inches(0.35), Inches(0.196)
    inset = Inches(0.13)
    tf = pill.text_frame
    tf.word_wrap = False
    tf.margin_left = int(inset + btn_w + Inches(0.09))
    tf.margin_right = Inches(0.07)
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = label
    r.font.name = "Calibri"
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = text_color
    pill.click_action.target_slide = target_slide
    btn = _add_jump_button(slide, target_slide,
                           left=int(left + inset),
                           top=int(top + (height - btn_h) / 2),
                           width=btn_w, height=btn_h, back=back)
    return pill, btn


# Post-work reference box (Teaching CLAUDE.md): the glyph is a fixed
# vocabulary — ✎ always means "a problem set", ▤ always "a teaching note" —
# and the default position is the bottom-RIGHT corner, overlapping the
# footer.  Route every call through these two constants (2026-08-27).
PS_GLYPH = "✎"


PS_BOX_XY = (Inches(10.17), Inches(6.53))


def _add_ps_pointer(slide, *, left=None, top=None, label="Problem Set 1",
                    width=Inches(2.5), height=Inches(0.5)):
    if left is None:
        left = PS_BOX_XY[0]
    if top is None:
        top = PS_BOX_XY[1]
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 int(left), int(top), int(width), int(height))
    try:
        shp.adjustments[0] = 0.28
    except Exception:
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = WHITE
    shp.line.color.rgb = GOLD
    shp.line.width = Pt(1.5)
    shp.shadow.inherit = False
    _add_drop_shadow(shp)
    tf = shp.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
    tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r1 = p.add_run(); r1.text = PS_GLYPH + "  "
    r1.font.name = "Calibri"; r1.font.size = Pt(16)
    r1.font.bold = True; r1.font.color.rgb = GOLD
    r2 = p.add_run(); r2.text = label
    r2.font.name = "Calibri"; r2.font.size = Pt(16)
    r2.font.bold = True; r2.font.color.rgb = NAVY
    return shp


RED_MW = RGBColor(0xD5, 0x1A, 0x1A)        # MW's ▼ red


TITLE_LOWER = {
    "a", "an", "the",
    "and", "but", "or", "nor", "for", "so", "yet",
    "as", "at", "by", "in", "into", "of", "off", "on", "onto", "out",
    "per", "to", "up", "via", "with", "from", "over", "than", "vs",
}


def _tc_word(word, force):
    """Capitalise one whitespace-delimited token, hyphen parts included."""
    parts = word.split("-")
    out = []
    for i, part in enumerate(parts):
        core = part.lstrip("“‘\"'([")
        lead = part[:len(part) - len(core)]
        stripped = core.rstrip("”’\"')]:,.?!")
        trail = core[len(stripped):]
        low = stripped.lower().rstrip(".")
        keep_low = (not force) and low in TITLE_LOWER
        if stripped and not keep_low and stripped[0].islower():
            stripped = stripped[0].upper() + stripped[1:]
        out.append(lead + stripped + trail)
        force = False          # only the first hyphen part can be forced
    return "-".join(out)


def _title_case(title):
    words = title.split(" ")
    idx = [i for i, w in enumerate(words) if w.strip()]
    if not idx:
        return title
    first, last = idx[0], idx[-1]
    out = []
    force_next = True
    for i, w in enumerate(words):
        if not w.strip():
            out.append(w)
            continue
        force = force_next or i == first or i == last
        out.append(_tc_word(w, force))
        force_next = w.rstrip().endswith((":", "?", "!", "—", "–"))
    return " ".join(out)


# 2026-08-25 (Nico, measured off the final video decks): an item with no
# description shown renders its single line at the TOP of the reserved
# two-row box, which sits high against the gold circle.  Every SHADED item
# is nudged down by exactly this much to centre it; the current topic,
# which fills its box with title + description, does not move.
DIM_DROP = 85064                           # 0.093 in


# 2026-08-30 (Nico): LR is an index like a digit — it means "long run",
# so P_LR and Q_LR take a real subscript wherever they appear.
SYMBOL_RE = re.compile(
    r"(MC|MR|SMC|EMC|AVC|ATC|LAC|LMC|TC|[PQDSwL])([\u2032\u2019']?)"
    r"(?:_(\w+)|([0-9]|Peak|LR|Labor|min|max))(?![A-Za-z0-9])")


SUBSCRIPT_BASELINE = "-25000"      # what PowerPoint writes for subscript


def _split_symbol_runs(para):
    """Rewrite the runs of one paragraph so every P/Q/D/S symbol carries an
    italic base and a subscript index."""
    changed = False
    for run in list(para.runs):
        text = run.text or ""
        if not SYMBOL_RE.search(text):
            continue
        pieces = []            # (text, is_subscript)
        pos = 0
        for m in SYMBOL_RE.finditer(text):
            if m.start() > pos:
                pieces.append((text[pos:m.start()], False))
            pieces.append((m.group(1) + m.group(2), False))
            pieces.append((m.group(3) or m.group(4), True))
            pos = m.end()
        if pos < len(text):
            pieces.append((text[pos:], False))
        if len(pieces) < 2:
            continue
        r_el = run._r
        rPr = r_el.find(qn('a:rPr'))
        anchor = r_el
        for i, (txt, is_sub) in enumerate(pieces):
            if i == 0:
                run.text = txt
                if rPr is not None:
                    rPr.set('i', '1')
                continue
            new_r = copy.deepcopy(r_el)
            for t_el in new_r.findall(qn('a:t')):
                t_el.text = txt
            npr = new_r.find(qn('a:rPr'))
            if npr is not None:
                npr.set('i', '1')
                if is_sub:
                    npr.set('baseline', SUBSCRIPT_BASELINE)
                else:
                    npr.attrib.pop('baseline', None)
            anchor.addnext(new_r)
            anchor = new_r
        changed = True
    return changed


def apply_symbol_subscripts(prs):
    """Deck-wide: italic base + subscript index for every P/Q/D/S symbol."""
    n = 0
    for slide in prs.slides:
        for shp in slide.shapes:
            for tf in _iter_text_frames(shp):
                for para in tf.paragraphs:
                    if _split_symbol_runs(para):
                        n += 1
    return n


def _iter_text_frames(shp):
    if shp.has_text_frame:
        yield shp.text_frame
    if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shp.shapes:
            for tf in _iter_text_frames(child):
                yield tf
    if getattr(shp, "has_table", False) and shp.has_table:
        for row in shp.table.rows:
            for cell in row.cells:
                yield cell.text_frame

