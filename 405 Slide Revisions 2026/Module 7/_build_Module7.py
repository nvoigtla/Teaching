# ==========================================================================
#  _build_Module7.py — phase-1 scaffold for "Module 7 - Revised.pptx"
#
#  Oligopoly and Game Theory (87 slides, 1:1 with the source deck
#  "Module 7.pptx"; poll/video slides are positional stubs until the
#  phase-3 OOXML splice).
#
#  Helper layer (through the OMML/equation engine) copied VERBATIM from
#  Module 3/_build_Module3.py on 2026-07-29 — the proven chrome, box,
#  bullet, chart, and table primitives. M7-specific code starts at the
#  "MODULE 7" banner below.
# ==========================================================================

import copy
import math
import re
import shutil
import zipfile
from pathlib import Path

from lxml import etree as ET
from pptx import Presentation
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Cm, Inches, Pt

# Reuse all primitives from the template script (single source of truth).
from _build_template_samples import (
    FADED,
    FOOTER_TEXT,
    GOLD,
    GOLD_W,
    GRAY,
    MARGIN,
    MODULE_AGENDA,
    NAVY,
    RULE,
    RULE_W,
    SLIDE_H,
    SLIDE_W,
    WHITE,
    _add_bulleted_list,
    _add_rect,
    _add_text,
    _blank_slide,
    _draw_action_title,
    _set_bullet_char,
)

OUT_DIR = Path(__file__).parent


# --------------------------------------------------------------------------
# Title-case top bar – replaces the all-caps default from the template script.
# The user prefers academic-paper title case (each major word capitalized,
# small connectors like "and", "of", "for", "the" stay lowercase).
# --------------------------------------------------------------------------

def _draw_top_bar_tc(slide, section_tag):
    """Navy top bar with title-cased section tag (no uppercase forcing)."""
    bar_h = Inches(0.42)
    _add_rect(slide, 0, 0, SLIDE_W, bar_h, NAVY)
    _add_text(slide, MARGIN, 0, Inches(12), bar_h,
              section_tag, size=16, bold=True,
              color=WHITE, font="Calibri",
              anchor=MSO_ANCHOR.MIDDLE)


def _draw_footer(slide, footer_text, page_num):
    """Footer rule + gold accent + footer text/page number, with larger type
    than the template default so handout printouts remain legible."""
    _add_rect(slide, 0, Inches(7.15), SLIDE_W, Inches(0.02), RULE)
    _add_rect(slide, MARGIN, Inches(7.135), GOLD_W, Inches(0.05), GOLD)
    _add_text(slide, MARGIN, Inches(7.20), Inches(11), Inches(0.32),
              footer_text, size=12, color=GRAY)
    _add_text(slide, Inches(12.5), Inches(7.20), Inches(0.6), Inches(0.32),
              str(page_num), size=12, color=GRAY, align=PP_ALIGN.RIGHT)


# --------------------------------------------------------------------------
# Speaker-notes helper
# --------------------------------------------------------------------------

def _set_notes(slide, text):
    """Replace the slide's speaker notes with *text*."""
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.clear()
    notes_tf.text = text


# --------------------------------------------------------------------------
# Reusable shape primitives for diagrams
# --------------------------------------------------------------------------

def _add_filled_box(slide, left, top, width, height, label, *,
                    fill=NAVY, text_color=WHITE, line=None,
                    size=18, bold=True, font="Calibri"):
    """Filled rectangle with centered text."""
    left, top, width, height = int(left), int(top), int(width), int(height)
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height,
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = text_color
    return shp


def _add_outlined_box(slide, left, top, width, height, label, *,
                      line=NAVY, text_color=NAVY, fill=WHITE,
                      size=18, bold=True, line_w=1.25, font="Calibri",
                      rounded=False, shadow=False, corner_pct=0.06):
    """Outlined rectangle (white fill) with centered text.

    ``rounded=True`` switches the base shape to a rounded rectangle
    (corner-adjust = ``corner_pct``); ``shadow=True`` adds a soft drop
    shadow.  Both default off so existing flat call-sites are
    unaffected.
    """
    left, top, width, height = int(left), int(top), int(width), int(height)
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(
        shape_type, left, top, width, height,
    )
    if rounded:
        try: shp.adjustments[0] = corner_pct
        except Exception: pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    if shadow:
        _add_drop_shadow(shp)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = text_color
    return shp


def _add_convention_box(slide, left, top, width, height, *,
                          prefix=None, body=None, runs=None,
                          fill_rgb=None, border=None, line_w=1.0,
                          corner_pct=0.12, size=15, align=PP_ALIGN.LEFT,
                          font="Calibri", pad_h=None, pad_v=None,
                          line_spacing_pct=None):
    """Cream-fill / navy-border rounded-rect explanation callout.

    The "Convention" textbox pattern from slide 14 generalised — use it
    anywhere a slide needs a compact, visually-distinct box for a
    short conceptual explanation or notational convention.  Sits well
    below a table, beside a hero formula, or as a slide-wide footer.

    Two ways to populate the text:
      • ``prefix`` (bold) + ``body`` (regular) — simplest path; matches
        slide 14's "Convention:  <text>" pattern.
      • ``runs`` — a list of ``(text, {"bold": .., "italic": .., ...})``
        tuples for finer-grained styling (multi-line, mixed formatting).

    Style defaults follow the course-layer CLAUDE.md "Convention callout
    box" spec — cream fill, thin primary-color border, slight rounding,
    primary-color text.  Override ``fill_rgb`` / ``border`` only when you
    need a different accent.
    """
    fill = fill_rgb if fill_rgb is not None else RGBColor(0xFD, 0xF6, 0xE6)
    border = border if border is not None else NAVY

    left, top, width, height = int(left), int(top), int(width), int(height)
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border
    shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    try:
        shp.adjustments[0] = corner_pct
    except Exception:
        pass

    # Inset text box so the rounded corners breathe — matches slide 14.
    pad_h = Inches(0.20) if pad_h is None else pad_h
    pad_v = Inches(0.12) if pad_v is None else pad_v
    tb = slide.shapes.add_textbox(
        left + int(pad_h), top + int(pad_v),
        width - 2 * int(pad_h), height - 2 * int(pad_v),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0); tf.margin_bottom = Inches(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    def _style_run(r, opts):
        r.font.name = opts.get('font', font)
        r.font.size = Pt(opts.get('size', size))
        r.font.bold = opts.get('bold', False)
        r.font.italic = opts.get('italic', False)
        r.font.underline = opts.get('underline', False)
        r.font.color.rgb = opts.get('color', NAVY)

    if runs is not None:
        first = True
        for entry in runs:
            text, opts = entry if isinstance(entry, tuple) else (entry, {})
            if opts.get('newline') and not first:
                p = tf.add_paragraph()
            elif first:
                p = tf.paragraphs[0]
            else:
                # Same paragraph — append run to the most-recent paragraph.
                p = tf.paragraphs[-1]
            p.alignment = align
            r = p.add_run()
            r.text = text
            _style_run(r, opts)
            first = False
    else:
        p = tf.paragraphs[0]
        p.alignment = align
        if prefix:
            r1 = p.add_run(); r1.text = prefix
            _style_run(r1, {'bold': True, 'color': NAVY, 'size': size})
        if body:
            r2 = p.add_run(); r2.text = body
            _style_run(r2, {'color': NAVY, 'size': size})

    if line_spacing_pct is not None:
        for p_obj in tf.paragraphs:
            pPr = p_obj._p.get_or_add_pPr()
            for old in pPr.findall(qn('a:lnSpc')):
                pPr.remove(old)
            lnSpc = ET.Element(qn('a:lnSpc'))
            spcPct = ET.SubElement(lnSpc, qn('a:spcPct'))
            spcPct.set('val', str(int(line_spacing_pct * 1000)))
            pPr.insert(0, lnSpc)
    return shp


def _add_rounded_filled_box(slide, left, top, width, height, label, *,
                             fill=NAVY, text_color=WHITE, line=None,
                             size=18, bold=True, font="Calibri",
                             corner_pct=0.06, shadow=True):
    """Rounded-corner filled rectangle with centered text and soft drop shadow.

    Mirrors :func:`_add_filled_box` but renders ``MSO_SHAPE.ROUNDED_RECTANGLE``
    with the corner-adjust set to ``corner_pct`` (6 % per course CLAUDE.md
    "slight rounding") and a soft drop shadow via :func:`_add_drop_shadow`.
    """
    left, top, width, height = int(left), int(top), int(width), int(height)
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
    )
    try:
        shp.adjustments[0] = corner_pct
    except Exception:
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    if shadow:
        _add_drop_shadow(shp)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = text_color
    return shp


def _add_arrow(slide, start_xy, end_xy, *, color=NAVY, weight_pt=1.5,
               head=True, dash=None, head_size='med'):
    """Draw a line/arrow from start to end (in EMU/Inches values).

    EMU coordinates MUST be integers — PowerPoint rejects decimal values
    in <a:off>/<a:ext> and refuses to open the file. Cast to int defensively.

    ``dash`` accepts any OOXML preset-dash name (e.g., ``"dash"``,
    ``"dashDot"``, ``"sysDash"``).  Default ``None`` = solid line.

    ``head_size`` is the OOXML preset arrowhead size — one of ``'sm'``,
    ``'med'`` (default), or ``'lg'``.  Width and height are set
    together, so passing ``'lg'`` gives a noticeably larger tip while
    leaving the line weight unchanged.
    """
    sx, sy = int(start_xy[0]), int(start_xy[1])
    ex, ey = int(end_xy[0]), int(end_xy[1])
    line = slide.shapes.add_connector(1, sx, sy, ex, ey)  # 1 = STRAIGHT
    line.line.color.rgb = color
    line.line.width = Pt(weight_pt)
    ln = line.line._get_or_add_ln()
    if dash is not None:
        for old in ln.findall(qn('a:prstDash')):
            ln.remove(old)
        prst = ET.SubElement(ln, qn('a:prstDash'))
        prst.set('val', dash)
    if head:
        tailEnd = ET.SubElement(ln, qn('a:tailEnd'))
        tailEnd.set('type', 'triangle')
        tailEnd.set('w', head_size)
        tailEnd.set('h', head_size)
    return line


def _add_wavy_line(slide, x_start, x_end, y_center, *,
                    amplitude=None, cycles=1.75, segments=36,
                    color=NAVY, weight_pt=1.5):
    """Horizontal sinusoidal line from x_start to x_end at y_center.

    Renders a polyline approximation of ``sin(2π · t · cycles)`` inside a
    custGeom shape so the line reads as a gentle wave rather than a
    straight connector.  ``amplitude`` is the peak-to-baseline height in
    EMU; defaults to ~0.04".  ``segments`` controls how finely the wave
    is discretised — 30+ is smooth enough that the polyline reads as a
    curve.  No arrowhead.
    """
    if amplitude is None:
        amplitude = Inches(0.04)
    L = int(x_end - x_start)
    A = int(amplitude)
    if L == 0:
        return None
    flip = "1" if L < 0 else "0"
    bbox_left = int(min(x_start, x_end))
    bbox_top = int(y_center - A)
    bbox_w = abs(L)
    bbox_h = 2 * A

    pts = []
    for i in range(segments + 1):
        t = i / segments
        lx = int(round(t * 100000))
        sin_val = math.sin(2 * math.pi * t * cycles)
        # Path coords: y=0 is top.  Centre at 50000; +1·amplitude → top
        # (0), −1·amplitude → bottom (100000).
        ly = int(round(50000 - sin_val * 50000))
        pts.append((lx, ly))

    path_segs = [f'<a:moveTo><a:pt x="{pts[0][0]}" y="{pts[0][1]}"/></a:moveTo>']
    for lx, ly in pts[1:]:
        path_segs.append(f'<a:lnTo><a:pt x="{lx}" y="{ly}"/></a:lnTo>')
    path_inner = ''.join(path_segs)

    color_hex = f'{color[0]:02X}{color[1]:02X}{color[2]:02X}'
    width_emu = int(weight_pt * 12700)

    P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    A_NS_LOCAL = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    sp_xml = (
        f'<p:sp xmlns:p="{P_NS}" xmlns:a="{A_NS_LOCAL}">'
        f'<p:nvSpPr><p:cNvPr id="0" name="WavyLine"/>'
        f'<p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr>'
        f'<a:xfrm flipH="{flip}">'
        f'<a:off x="{bbox_left}" y="{bbox_top}"/>'
        f'<a:ext cx="{bbox_w}" cy="{bbox_h}"/>'
        f'</a:xfrm>'
        f'<a:custGeom>'
        f'<a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/>'
        f'<a:rect l="0" t="0" r="0" b="0"/>'
        f'<a:pathLst>'
        f'<a:path w="100000" h="100000" fill="none">'
        f'{path_inner}'
        f'</a:path>'
        f'</a:pathLst>'
        f'</a:custGeom>'
        f'<a:noFill/>'
        f'<a:ln w="{width_emu}" cap="rnd">'
        f'<a:solidFill><a:srgbClr val="{color_hex}"/></a:solidFill>'
        f'</a:ln>'
        f'</p:spPr>'
        f'</p:sp>'
    )
    elem = ET.fromstring(sp_xml)
    slide.shapes._spTree.append(elem)
    return elem


def _add_arrow_shape(slide, left, top, width, height, *,
                     direction="right", fill=GOLD, line=None):
    """Block arrow shape (the 'we-are-here' indicator).

    direction: "right" (default), "left", "up", or "down".
    """
    left, top, width, height = int(left), int(top), int(width), int(height)
    geom_map = {
        "left": MSO_SHAPE.LEFT_ARROW,
        "right": MSO_SHAPE.RIGHT_ARROW,
        "up": MSO_SHAPE.UP_ARROW,
        "down": MSO_SHAPE.DOWN_ARROW,
    }
    geom = geom_map.get(direction, MSO_SHAPE.RIGHT_ARROW)
    shp = slide.shapes.add_shape(geom, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp



# --------------------------------------------------------------------------
# Layout 2 — Section Header / Agenda (parameterized)
# --------------------------------------------------------------------------

def make_section_agenda(prs, page_num, *, current_part_idx=None,
                        current_sub_idx=None,
                        section_tag="Module 3 · Agenda",
                        title="Agenda"):
    """Render an agenda / section-divider slide.

    ``current_part_idx`` makes that Part navy and fades the others.
    ``current_sub_idx`` (optional, only meaningful with current_part_idx)
    further restricts the navy highlight inside the current Part to a
    single sub-bullet — the other subs render in faded gray, alongside
    Part 2.  Used by intra-Part dividers (e.g., the new Short Run agenda
    on page 13, or the Long Run agenda on page 31).
    """
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, section_tag)
    _draw_action_title(slide, title)

    y = Inches(1.85)
    part_title_h = Inches(0.6)
    sub_list_h = Inches(1.35)
    block_gap = Inches(0.15)

    for idx, part in enumerate(MODULE_AGENDA):
        if current_part_idx is None:
            color = NAVY  # preview: highlight all Parts
        else:
            color = NAVY if idx == current_part_idx else FADED

        # 2026-05-19: per-sub fade.  When this is the current Part AND a
        # specific sub-index is highlighted, give every other sub the
        # FADED color so the deck shows "you are HERE" within the Part.
        if (current_part_idx is not None and idx == current_part_idx
                and current_sub_idx is not None):
            sub_colors = [
                NAVY if i == current_sub_idx else FADED
                for i in range(len(part["subs"]))
            ]
        else:
            sub_colors = None

        _add_text(slide, MARGIN, y, RULE_W, part_title_h,
                  part["title"], size=30, bold=True, color=color,
                  font="Calibri")
        y += part_title_h

        _add_bulleted_list(
            slide,
            left=MARGIN + Inches(0.4),
            top=y,
            width=RULE_W - Inches(0.4),
            height=sub_list_h,
            items=part["subs"],
            size=24, color=color, bullet_color=color,
            # 2026-05-19: sub-bullets tightened from 10 → 3 pt to match
            # the deck-wide rhythm (sub-bullets cluster under their
            # parent rather than breathing out).  Applied to both
            # Parts on every agenda divider in the deck.
            line_spacing_pts=3,
            autonum_scheme='alphaLcPeriod',
            colors=sub_colors,
        )
        y += sub_list_h + block_gap

    _draw_footer(slide, FOOTER_TEXT, page_num)
    return slide


# --------------------------------------------------------------------------
# Layout 3 — Content bulleted
# --------------------------------------------------------------------------

def make_content_bulleted(prs, page_num, section_tag, title, bullets, *,
                          size=24, sub_size=None, line_spacing_pts=18,
                          sub_line_spacing_pts=None,
                          extras=None, bullets_top=None):
    """bullets: list of (text, level) tuples OR plain strings (level=0).

    ``bullets_top`` overrides the default body-region start (Inches(1.85))
    — useful when the slide also hosts large diagrams below the bullets
    and the bullets need to lift up to avoid overlap.

    ``sub_line_spacing_pts`` overrides the legacy
    ``max(6, line_spacing_pts - 8)`` formula for sub-bullet space-before.
    """
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, section_tag)
    _draw_action_title(slide, title)

    normalized = [(b, 0) if isinstance(b, str) else b for b in bullets]

    if bullets_top is None:
        bullets_top = Inches(1.85)

    _add_hierarchical_bullets(
        slide,
        left=MARGIN,
        top=bullets_top,
        width=RULE_W,
        height=Inches(5.0),
        items=normalized,
        size=size,
        sub_size=sub_size,
        line_spacing_pts=line_spacing_pts,
        sub_line_spacing_pts=sub_line_spacing_pts,
    )

    if extras is not None:
        extras(slide)

    _draw_footer(slide, FOOTER_TEXT, page_num)
    return slide


def _inject_bullet_lst_style(tf, *, size, sub_size,
                              main_color=NAVY, sub_color=GRAY,
                              main_char='▪', sub_char='–',
                              main_space_pts=None, sub_space_pts=None):
    """Inject <a:lstStyle> into a text frame defining per-level defaults.

    Defines lvl1pPr (main bullets, lvl="0") and lvl2pPr (sub bullets,
    lvl="1") with bullet character, indent, color, font, and space-
    before defaults.  Enables PowerPoint's native Tab / Shift+Tab to
    auto-reformat a bullet when its level changes: runs / paragraphs
    that do NOT explicitly override these properties inherit them.

    main_space_pts / sub_space_pts: if set, the <a:spcBef> for each
    level.  Sub-bullets default to a tighter spacing than main bullets
    (3 pt vs ~10 pt) per the deck's pedagogical preference: sub-bullets
    visually cluster under their parent, main bullets give the eye a
    larger break.
    """
    txBody = tf._txBody
    # Remove any prior lstStyle
    for old in txBody.findall(qn('a:lstStyle')):
        txBody.remove(old)
    lstStyle = ET.Element(qn('a:lstStyle'))
    levels = [
        ('a:lvl1pPr', 342900, -342900, main_color, main_char, size, main_space_pts),
        ('a:lvl2pPr', 685800, -228600, sub_color, sub_char, sub_size, sub_space_pts),
    ]
    for tag, mar_l, indent, color, char, sz, space_pts in levels:
        lvl = ET.SubElement(lstStyle, qn(tag))
        lvl.set('marL', str(mar_l))
        lvl.set('indent', str(indent))
        # spcBef must precede bullet attributes (OOXML schema order).
        if space_pts is not None:
            spcBef = ET.SubElement(lvl, qn('a:spcBef'))
            spcPts = ET.SubElement(spcBef, qn('a:spcPts'))
            spcPts.set('val', str(int(space_pts * 100)))
        bc = ET.SubElement(lvl, qn('a:buClr'))
        sc = ET.SubElement(bc, qn('a:srgbClr'))
        sc.set('val', '{:02X}{:02X}{:02X}'.format(color[0], color[1], color[2]))
        bf = ET.SubElement(lvl, qn('a:buFont'))
        bf.set('typeface', 'Calibri')
        bch = ET.SubElement(lvl, qn('a:buChar'))
        bch.set('char', char)
        drp = ET.SubElement(lvl, qn('a:defRPr'))
        drp.set('sz', str(int(sz * 100)))
        drp.set('b', '0')
        sf = ET.SubElement(drp, qn('a:solidFill'))
        sfc = ET.SubElement(sf, qn('a:srgbClr'))
        sfc.set('val', '{:02X}{:02X}{:02X}'.format(color[0], color[1], color[2]))
        lt = ET.SubElement(drp, qn('a:latin'))
        lt.set('typeface', 'Calibri')
    # Insert lstStyle after bodyPr (proper element order in <a:txBody>)
    bodyPr = txBody.find(qn('a:bodyPr'))
    if bodyPr is not None:
        bodyPr.addnext(lstStyle)
    else:
        txBody.insert(0, lstStyle)


def _add_hierarchical_bullets(slide, left, top, width, height, items,
                              *, size=24, sub_size=None, line_spacing_pts=18,
                              sub_line_spacing_pts=None):
    """Render bullets with indent levels.

    Bullet item forms:
        (text, level)                       — simple, defaults from level
        (text, level, opts)                 — opts dict overrides
        (runs_list, level, opts)            — multi-run paragraph

    text:
        - str  → a single run with text
        - ''   → empty paragraph (visual spacer; no run)
        - list → multi-run: list of ``(run_text, run_opts)`` tuples

    Paragraph-level opts (all optional):
        bullet_style: 'main' (▪ NAVY) | 'sub' (– GRAY) | 'arrow' (no bullet,
            plain left-indent; the run text supplies the leader char like
            "→" or Wingdings) | 'none'
        mar_l, indent: bullet positioning (EMU)
        space_before_pts: spcBef in pts (overrides legacy formula)
        size, color, bold, italic: defaults applied to every run unless
            the run_opts override them.

    Run-level opts (run_opts in a runs_list tuple):
        font_name (default 'Calibri'), size, color, bold, italic,
        underline (bool), wingdings (bool — emits <a:sym typeface="Wingdings"/>
        so a private-use-area character renders as its Wingdings glyph).

    Each paragraph also receives a ``lvl="N"`` attribute when level > 0
    so PowerPoint's Tab / Shift-Tab outline navigation can find an
    explicit outline level.
    """
    if sub_size is None:
        sub_size = size - 4

    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    # 2026-05-18: inject <a:lstStyle> with per-level defaults so that
    # PowerPoint's Tab / Shift+Tab keyboard shortcuts auto-reformat the
    # bullet when its outline level changes.  Default 'main'/'sub'
    # bullets (no custom indent, level < 2) skip explicit pPr/run
    # styling — they inherit from lstStyle.  Bullets with custom indent
    # OR level >= 2 OR explicit size/color overrides still apply
    # explicit styling.
    # 2026-05-18 (later): also moved space-before into lstStyle so a
    # paragraph demoted via Tab picks up the sub-bullet's tighter spcBef
    # automatically.  Sub-bullet default: 3 pt (was 6 pt) — per user
    # preference, sub-bullets cluster tightly under their parent.
    sub_default_space = (sub_line_spacing_pts if sub_line_spacing_pts is not None
                          else 3)
    _inject_bullet_lst_style(tf, size=size, sub_size=sub_size,
                              main_space_pts=line_spacing_pts,
                              sub_space_pts=sub_default_space)

    for i, item in enumerate(items):
        if len(item) == 2:
            text, level = item
            opts = {}
        else:
            text, level, opts = item
            opts = opts or {}

        # Determine inheritance up-front so spcBef logic below can use it.
        style = opts.get('bullet_style', 'main' if level == 0 else 'sub')
        has_custom_indent = 'mar_l' in opts or 'indent' in opts
        inherits_from_lst = (style in ('main', 'sub')
                              and not has_custom_indent
                              and level < 2)

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        pPr = p._p.get_or_add_pPr()

        # Space-before.  When the paragraph inherits from lstStyle and the
        # caller has not overridden, skip the per-pPr spcBef so that Tab-
        # induced level changes pick up the new level's spcBef from
        # lstStyle.  Otherwise set explicitly.
        if i > 0:
            sp_override = opts.get('space_before_pts')
            if sp_override is not None:
                spcBef = ET.SubElement(pPr, qn('a:spcBef'))
                pts = ET.SubElement(spcBef, qn('a:spcPts'))
                pts.set('val', str(sp_override * 100))
            elif not inherits_from_lst:
                if level == 0:
                    sp = line_spacing_pts
                else:
                    sp = sub_default_space
                spcBef = ET.SubElement(pPr, qn('a:spcBef'))
                pts = ET.SubElement(spcBef, qn('a:spcPts'))
                pts.set('val', str(sp * 100))
            # else: inherit from lstStyle.

        # Outline-level attribute (enables PowerPoint Tab/Shift-Tab)
        if level > 0:
            pPr.set('lvl', str(level))

        # Bullet styling — default 'main' (lvl 0) and 'sub' (lvl 1)
        # inherit from lstStyle when no custom indent is requested.
        # 'arrow' and 'none' suppress the lstStyle bullet via <a:buNone/>.
        # Level >= 2 still uses explicit _set_bullet_char (lstStyle here
        # only defines lvl1pPr and lvl2pPr).
        if style == 'main':
            if not inherits_from_lst:
                _set_bullet_char(p, char='▪', color=NAVY,
                                  mar_l=opts.get('mar_l', 342900),
                                  indent=opts.get('indent', -342900),
                                  size_pct=100)
        elif style == 'sub':
            if not inherits_from_lst:
                default_mar = 342900 + level * 342900
                _set_bullet_char(p, char='–', color=GRAY,
                                  mar_l=opts.get('mar_l', default_mar),
                                  indent=opts.get('indent', -228600),
                                  size_pct=100)
        elif style == 'arrow':
            # Plain left-indent, no bullet glyph; user text supplies leader.
            pPr.set('marL', str(opts.get('mar_l', 457200)))
            if 'indent' in opts:
                pPr.set('indent', str(opts['indent']))
            # Explicit <a:buNone/> so the lstStyle bullet doesn't bleed
            # through.
            ET.SubElement(pPr, qn('a:buNone'))
        elif style == 'none':
            ET.SubElement(pPr, qn('a:buNone'))

        # Empty paragraph (spacer) — no run; suppress the bullet so the
        # empty line doesn't render an orphan glyph.
        if text == '':
            if inherits_from_lst:
                ET.SubElement(pPr, qn('a:buNone'))
            continue

        # Normalize text to a runs list
        if isinstance(text, str):
            runs = [(text, {})]
        else:
            runs = text  # already a list of (run_text, run_opts) tuples

        # Paragraph-level explicit overrides (None = inherit from lstStyle).
        para_size_override = opts.get('size')
        para_color_override = opts.get('color')
        para_bold_override = opts.get('bold')
        para_italic_override = opts.get('italic')

        for run_text, run_opts in runs:
            run_opts = run_opts or {}

            # Inline OMML math zone — append <a14:m><m:oMath>...</m:oMath></a14:m>
            # as a sibling of the regular <a:r> runs so the formula renders
            # in-line with surrounding prose.  ``run_text`` is interpreted
            # as raw OMML content (e.g., from `_formula_mp_ratio` or the
            # `_omml_*` builders).
            if run_opts.get('omml'):
                om_sz = run_opts.get('size')
                if om_sz is None:
                    om_sz = para_size_override
                if om_sz is None:
                    om_sz = size if level == 0 else sub_size
                om_clr = run_opts.get('color') or para_color_override
                if om_clr is None:
                    om_clr = NAVY if level == 0 else GRAY
                sz_centi = int(om_sz * 100)
                clr_hex = '{:02X}{:02X}{:02X}'.format(
                    om_clr[0], om_clr[1], om_clr[2])
                a14_xml = (
                    f'<a14:m xmlns:a14="{A14_NS}" '
                    f'xmlns:m="{M_NS}" xmlns:a="{A_NS}">'
                    f'<m:oMath>{run_text}</m:oMath>'
                    f'</a14:m>'
                )
                a14_elem = ET.fromstring(a14_xml)
                for r_elem in a14_elem.iter(qn('m:r')):
                    arPr = r_elem.find(qn('a:rPr'))
                    if arPr is None:
                        arPr = ET.Element(qn('a:rPr'))
                        r_elem.insert(0, arPr)
                    arPr.set('sz', str(sz_centi))
                    if arPr.get('lang') is None:
                        arPr.set('lang', 'en-US')
                    for sf in arPr.findall(qn('a:solidFill')):
                        arPr.remove(sf)
                    # fill must come BEFORE a:latin (schema order), else
                    # PowerPoint ignores the color
                    sf = arPr.makeelement(qn('a:solidFill'), {})
                    srgb = ET.SubElement(sf, qn('a:srgbClr'))
                    srgb.set('val', clr_hex)
                    arPr.insert(0, sf)
                p._p.append(a14_elem)
                continue

            run = p.add_run()
            run.text = run_text
            run.font.name = run_opts.get('font_name', 'Calibri')

            # Size: run-opt > para-opt > inherit (for inherits_from_lst
            # cases) or fall back to level default (for explicit-styled
            # paragraphs).
            run_sz = run_opts.get('size')
            sz = run_sz if run_sz is not None else para_size_override
            if sz is None and not inherits_from_lst:
                sz = size if level == 0 else sub_size
            if sz is not None:
                run.font.size = Pt(sz)

            run_clr = run_opts.get('color')
            clr = run_clr if run_clr is not None else para_color_override
            if clr is None and not inherits_from_lst:
                clr = NAVY if level == 0 else GRAY
            if clr is not None:
                run.font.color.rgb = clr

            # Bold: only set if explicitly requested.  Default = inherit
            # (lstStyle defRPr has b="0").
            run_b = run_opts.get('bold')
            b = run_b if run_b is not None else para_bold_override
            if b is not None:
                run.font.bold = b

            it = run_opts.get('italic', para_italic_override)
            if it is not None:
                run.font.italic = it
            if run_opts.get('underline'):
                run.font.underline = True
            if run_opts.get('wingdings'):
                # Add <a:sym typeface="Wingdings"/> so private-use-area
                # characters render as their Wingdings glyphs.
                rPr = run._r.find(qn('a:rPr'))
                if rPr is None:
                    rPr = run._r.makeelement(qn('a:rPr'), {})
                    run._r.insert(0, rPr)
                sym = ET.SubElement(rPr, qn('a:sym'))
                sym.set('typeface', 'Wingdings')

    return box


# --------------------------------------------------------------------------
# Layout-3 variant — diagram canvas (action title + free-form shapes below).
# Used for the agenda flowchart and the Big Picture diagram so they live
# inside the same visual chrome as content slides but render a diagram
# instead of bullets.
# --------------------------------------------------------------------------

def make_diagram_slide(prs, page_num, section_tag, title, draw_diagram):
    """Action-title slide with free-form diagram region below.

    draw_diagram: callable(slide) that renders the diagram in the body
    region (approximately y = 1.85 to 6.95).
    """
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, section_tag)
    _draw_action_title(slide, title)
    draw_diagram(slide)
    _draw_footer(slide, FOOTER_TEXT, page_num)
    return slide


# --------------------------------------------------------------------------
# Layout 5 — Poll slide (A./B./C./D. options + POLL pill, no QR box)
# --------------------------------------------------------------------------

def _draw_poll_pill(slide, *, position='top-right',
                     fill=NAVY, text_color=WHITE, dot_color=GOLD,
                     width=None, height=None, text_size=14,
                     shadow=False):
    """Small 'POLL' pill, right-aligned.

    position: 'top-right' (default, just below the top bar) or
              'bottom-right' (bottom band aligned with discussion break).
    width / height: override the default 1.05" × 0.42" pill size.
    text_size: "POLL" font size (pt).  Scales up when the pill is enlarged.
    dot_color: pass None to skip the leading accent dot.
    shadow: add a soft drop shadow to the pill (matches discussion-break
        chrome).  Default False to preserve the original flat top-right
        pill look.
    """
    pill_w = width if width is not None else Inches(1.05)
    pill_h = height if height is not None else Inches(0.42)
    pill_x = SLIDE_W - MARGIN - pill_w
    if position == 'bottom-right':
        # Bottom-aligned with where _add_discussion_break ends
        # (top=6.25, height=0.72 → bottom=6.97).  Pill height varies, so
        # pick top so the pill BOTTOM lands on the same 6.97 baseline.
        pill_y = Inches(6.97) - pill_h
    else:
        pill_y = Inches(0.55)

    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  pill_x, pill_y, pill_w, pill_h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    if shadow:
        _add_drop_shadow(shp)
    else:
        shp.shadow.inherit = False
    try:
        shp.adjustments[0] = 0.5
    except Exception:
        pass
    _add_text(slide, pill_x, pill_y, pill_w, pill_h,
              "POLL", size=text_size, bold=True, color=text_color, font="Calibri",
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if dot_color is not None:
        dot_size = Inches(0.16)
        dot_x = pill_x - Inches(0.16) - dot_size
        dot_y = pill_y + (pill_h - dot_size) // 2
        _add_rect(slide, dot_x, dot_y, dot_size, dot_size, dot_color)


def make_poll_slide(prs, page_num, section_tag, title, options, *,
                    instructions="Respond at PollEv.com/nvoigtlaender",
                    size=30, line_spacing_pts=22):
    """Layout 5 – Poll slide.

    options: list of strings (A./B./C./D. auto-numbering applied).
    """
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, section_tag)
    _draw_action_title(slide, title)
    _draw_poll_pill(slide)

    _add_bulleted_list(
        slide,
        left=MARGIN,
        top=Inches(2.0),
        width=RULE_W,
        height=Inches(4.4),
        items=options,
        size=size, color=NAVY, bullet_color=NAVY,
        line_spacing_pts=line_spacing_pts,
        autonum_scheme='alphaUcPeriod',
    )

    _add_text(slide, MARGIN, Inches(6.55), RULE_W, Inches(0.4),
              instructions, size=16, italic=True, color=GRAY,
              align=PP_ALIGN.RIGHT)

    _draw_footer(slide, FOOTER_TEXT, page_num)
    return slide


# --------------------------------------------------------------------------
# Slide-jump hyperlinks — make a run or shape clickable in slideshow
# mode so it advances PowerPoint to a different slide.  python-pptx
# doesn't expose this directly, so we manipulate the rPr / cNvPr XML
# and register the slide relationship via the public OPC API.
# --------------------------------------------------------------------------

def _add_slide_jump_hyperlink_run(source_slide, run, target_slide,
                                    *, lock_color=True, underline=True):
    """Make ``run`` a clickable hyperlink that advances to ``target_slide``.

    ``lock_color=True`` adds the Office hyperlinkcolor extension so the
    run's explicit ``solidFill`` color is preserved at render time —
    without it, PowerPoint repaints hyperlink text in the theme's
    ``hlink`` color (a light blue) regardless of what the rPr says.

    ``underline=True`` sets ``u="sng"`` so the hyperlink stays visibly
    underlined even when the color-lock suppresses the theme styling.
    """
    rId = source_slide.part.relate_to(target_slide.part, RT.SLIDE)
    rPr = run._r.get_or_add_rPr()
    if underline:
        rPr.set('u', 'sng')
    for hl in rPr.findall(qn('a:hlinkClick')):
        rPr.remove(hl)
    hlinkClick = ET.SubElement(rPr, qn('a:hlinkClick'))
    hlinkClick.set(qn('r:id'), rId)
    hlinkClick.set('action', 'ppaction://hlinksldjump')
    if lock_color:
        AHYP_NS = 'http://schemas.microsoft.com/office/drawing/2018/hyperlinkcolor'
        ext_xml = (
            f'<a:extLst xmlns:a="{A_NS}">'
            f'<a:ext uri="{{A12FA001-AC4F-418D-AE19-62706E023703}}">'
            f'<ahyp:hlinkClr xmlns:ahyp="{AHYP_NS}" val="tx"/>'
            f'</a:ext>'
            f'</a:extLst>'
        )
        hlinkClick.append(ET.fromstring(ext_xml))


def _add_slide_jump_hyperlink_shape(source_slide, shape, target_slide):
    """Make ``shape`` clickable so the whole shape jumps to ``target_slide``."""
    rId = source_slide.part.relate_to(target_slide.part, RT.SLIDE)
    nvSpPr = shape._element.find(qn('p:nvSpPr'))
    cNvPr = nvSpPr.find(qn('p:cNvPr')) if nvSpPr is not None else None
    if cNvPr is None:
        return
    for hl in cNvPr.findall(qn('a:hlinkClick')):
        cNvPr.remove(hl)
    hlinkClick = ET.SubElement(cNvPr, qn('a:hlinkClick'))
    hlinkClick.set(qn('r:id'), rId)
    hlinkClick.set('action', 'ppaction://hlinksldjump')


# --------------------------------------------------------------------------
# Image helpers — embed source-deck images into the new deck.
# Images are pre-extracted to _source_images/slide{N}_{rId}.{ext}.
# --------------------------------------------------------------------------

SRC_IMG_DIR = Path(__file__).parent / "_source_images"


def _add_drop_shadow(shape, *, blur="50800", dist="38100",
                      direction="2700000", alpha="45000"):
    """Add a soft drop shadow to any shape that exposes spPr (pictures,
    rectangles, rounded rects).  Default: 4pt blur, 3pt offset, 45° down-
    right, 45% opacity black.  Used deck-wide for figures/boxes."""
    try:
        spPr = shape._element.spPr
    except AttributeError:
        # Fallback for shapes whose XML element exposes spPr only via find()
        spPr = shape._element.find(qn('p:spPr'))
    if spPr is None:
        return shape
    for old in spPr.findall(qn('a:effectLst')):
        spPr.remove(old)
    effLst = ET.SubElement(spPr, qn('a:effectLst'))
    outerShdw = ET.SubElement(effLst, qn('a:outerShdw'))
    outerShdw.set('blurRad', str(blur))
    outerShdw.set('dist', str(dist))
    outerShdw.set('dir', str(direction))
    outerShdw.set('algn', 'tl')
    outerShdw.set('rotWithShape', '0')
    rgb = ET.SubElement(outerShdw, qn('a:srgbClr'))
    rgb.set('val', '000000')
    a = ET.SubElement(rgb, qn('a:alpha'))
    a.set('val', str(alpha))
    return shape


def _add_graphicframe_shadow(slide, left, top, width, height, *,
                              shadow_alpha=45000):
    """White backing rectangle with an outerShdw effect, behind a
    graphicFrame (table or chart).  graphicFrames can't host
    a:effectLst directly; this rect supplies the shadow projected
    OUTSIDE its bounds.

    Charts have transparent plot areas by default, so a coloured backing
    bleeds through and tints the figure.  Use white instead — the chart
    sees white through its own transparent areas (clean background) and
    the shadow renders only at the visible edges.  Call BEFORE adding
    the table/chart so z-order is correct.
    """
    shdw = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        int(left), int(top), int(width), int(height),
    )
    shdw.fill.solid()
    shdw.fill.fore_color.rgb = WHITE
    shdw.line.fill.background()
    shdw.shadow.inherit = False
    sp_pr = shdw._element.spPr
    # Strip any default <a:effectLst/> python-pptx may have inserted
    # (duplicate effectLst makes PowerPoint refuse to open the file).
    for old in sp_pr.findall(qn('a:effectLst')):
        sp_pr.remove(old)
    effLst = ET.SubElement(sp_pr, qn('a:effectLst'))
    outerShdw = ET.SubElement(effLst, qn('a:outerShdw'))
    outerShdw.set('blurRad', '50800')
    outerShdw.set('dist', '38100')
    outerShdw.set('dir', '2700000')
    outerShdw.set('algn', 'tl')
    outerShdw.set('rotWithShape', '0')
    rgb = ET.SubElement(outerShdw, qn('a:srgbClr'))
    rgb.set('val', '000000')
    a = ET.SubElement(rgb, qn('a:alpha'))
    a.set('val', str(int(shadow_alpha)))
    return shdw


def _add_source_image(slide, src_slide_no, rid, *, left, top, width=None,
                      height=None, shadow=True, rounded=False):
    """Place a source-deck image on the new slide.

    `shadow=True` (default) adds a soft drop shadow so figures pop off the
    background — applied deck-wide per the latest visual direction.  Set
    `shadow=False` for niche cases (transparent PNGs, screenshots that
    already include a shadow, etc.).

    `rounded=True` gives real photographs the deck-standard rounded corners
    (routes through :func:`_apply_picture_style`, which sets both the
    rounded geometry and the drop shadow).  Leave False for logos,
    screenshots, and framed images per the Pictures guideline.
    """
    candidates = list(SRC_IMG_DIR.glob(f"slide{src_slide_no}_{rid}.*"))
    if not candidates:
        return None
    img = candidates[0]
    kwargs = {"left": int(left), "top": int(top)}
    if width is not None:
        kwargs["width"] = int(width)
    if height is not None:
        kwargs["height"] = int(height)
    pic = slide.shapes.add_picture(str(img), **kwargs)
    if rounded:
        _apply_picture_style(pic, corner_pct=6)
    elif shadow:
        _add_drop_shadow(pic)
    return pic


def _apply_picture_style(pic, *, corner_pct=8,
                          shadow_blur=50800, shadow_dist=38100,
                          shadow_dir=2700000, shadow_alpha=50000):
    """Apply rounded corners + drop shadow to a picture shape.

    corner_pct: rounded-corner radius as percent of shorter side (8 ≈ subtle).
    shadow_blur/dist: EMU; defaults give a soft 4pt blur, 3pt offset.
    shadow_dir: 2700000 = 45° down-right (standard).
    shadow_alpha: 50000 = 50% opacity black shadow.
    """
    spPr = pic._element.find(qn('p:spPr'))
    if spPr is None:
        return pic
    # Replace any existing prstGeom with roundRect at corner_pct
    for old in spPr.findall(qn('a:prstGeom')):
        spPr.remove(old)
    prstGeom = ET.Element(qn('a:prstGeom'))
    prstGeom.set('prst', 'roundRect')
    avLst = ET.SubElement(prstGeom, qn('a:avLst'))
    gd = ET.SubElement(avLst, qn('a:gd'))
    gd.set('name', 'adj')
    gd.set('fmla', f'val {int(corner_pct * 1000)}')
    # prstGeom must come after a:xfrm
    xfrm = spPr.find(qn('a:xfrm'))
    if xfrm is not None:
        xfrm.addnext(prstGeom)
    else:
        spPr.insert(0, prstGeom)
    # Replace any existing effectLst with a single outer shadow
    for old in spPr.findall(qn('a:effectLst')):
        spPr.remove(old)
    effectLst = ET.SubElement(spPr, qn('a:effectLst'))
    outerShdw = ET.SubElement(effectLst, qn('a:outerShdw'))
    outerShdw.set('blurRad', str(int(shadow_blur)))
    outerShdw.set('dist', str(int(shadow_dist)))
    outerShdw.set('dir', str(int(shadow_dir)))
    outerShdw.set('algn', 'tl')
    outerShdw.set('rotWithShape', '0')
    rgb = ET.SubElement(outerShdw, qn('a:srgbClr'))
    rgb.set('val', '000000')
    alpha = ET.SubElement(rgb, qn('a:alpha'))
    alpha.set('val', str(int(shadow_alpha)))
    return pic


# --------------------------------------------------------------------------
# Illustrative callout boxes — the "major-concept" badges, "Teaching Note"
# bars, and bottom-of-slide takeaway bands that recur throughout the source
# deck.  Replicating them faithfully (in template colors) is what makes the
# slides feel like the original, not a stripped-down rewrite.
# --------------------------------------------------------------------------

def _add_takeaway_bar(slide, text, *, top=Inches(6.4), width=None,
                       height=Inches(0.55), left=None,
                       fill=GOLD, text_color=WHITE,
                       size=20, font="Calibri", bold=True,
                       rounded=False, shadow=False):
    """Bottom-of-slide takeaway band — the 'major concept' callout.

    rounded=True  → ROUNDED_RECTANGLE with ~30% corner radius.
    shadow=True   → soft drop shadow (off by default to keep the flat
                    look on the majority of slides).
    left=None     → horizontally centered. Pass an EMU/Inches value to
                    pin the bar's left edge (used for hand-positioned
                    takeaways).
    """
    if width is None:
        width = Inches(9.6)
    if left is None:
        left = (SLIDE_W - width) // 2
    if not (rounded or shadow):
        return _add_filled_box(slide, left, top, width, height, text,
                                fill=fill, text_color=text_color,
                                size=size, bold=bold, font=font)
    # Inline build for the rounded / shadowed variant
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, int(left), int(top), int(width), int(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    if rounded:
        try: shape.adjustments[0] = 0.30
        except Exception: pass
    shape.shadow.inherit = False
    if shadow:
        _add_drop_shadow(shape)
    tf = shape.text_frame
    tf.margin_left = tf.margin_right = Inches(0.1)
    tf.margin_top = tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = text_color
    return shape


def _add_teaching_note(slide, text, *, top=Inches(6.6), width=None,
                        height=Inches(0.6), left=None,
                        rounded=False, pdf_icon=False, label_color=GOLD,
                        fill_rgb=None):
    """External-document reference card.

    Visually distinct from in-slide callouts: cream/parchment fill, navy
    DASHED border, italic navy text, with a small page-icon glyph on the
    left — clearly signals 'this links to an external Teaching Note doc'.

    Optional styling for the slide-32 (page 33) treatment:
      • ``left``       — pin to an absolute X (default: horizontally center)
      • ``rounded``    — rounded corners + soft drop shadow (default flat)
      • ``pdf_icon``   — render the page glyph as a white folded-corner
                        with bold "PDF" text, sized to nearly fill the card
                        (default: small navy folded-corner, no text)
      • ``label_color``— color of the "SEE TEACHING NOTE →" prefix
                        (default GOLD; pass NAVY for the page-33 look).
      • ``fill_rgb``   — override the cream/parchment card fill
                        (default RGBColor(0xF4, 0xF1, 0xEA)).
    """
    if width is None:
        width = Inches(8.0)
    if left is None:
        left = int((SLIDE_W - width) // 2)
    left, top, width, height = int(left), int(top), int(width), int(height)

    # Card: cream fill, dashed navy border.  ``rounded=True`` switches the
    # base shape to ROUNDED_RECTANGLE and adds a soft drop shadow.
    base_shape = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    card = slide.shapes.add_shape(base_shape, left, top, width, height)
    if rounded:
        try: card.adjustments[0] = 0.20
        except Exception: pass
    card.fill.solid()
    card.fill.fore_color.rgb = (fill_rgb if fill_rgb is not None
                                  else RGBColor(0xF4, 0xF1, 0xEA))  # parchment cream default
    card.line.color.rgb = NAVY
    card.line.width = Pt(1.25)
    card.shadow.inherit = False
    if rounded:
        _add_drop_shadow(card)
    # Dashed line style
    ln = card.line._get_or_add_ln()
    # Remove any existing prstDash
    for el in ln.findall(qn('a:prstDash')):
        ln.remove(el)
    prstDash = ET.SubElement(ln, qn('a:prstDash'))
    prstDash.set('val', 'dash')

    # Empty text on the card; we add the icon + text as separate textboxes
    tf = card.text_frame
    tf.text = ""

    # Small "page" icon on the left – a folded-corner shape.  In default
    # mode it's a small filled-navy glyph (logo-ish); in pdf_icon mode it's
    # a larger WHITE sheet with thin navy border + bold "PDF" text so it
    # reads as a generic PDF document at a glance.
    icon_size = Inches(0.5) if pdf_icon else Inches(0.4)
    icon_x = left + Inches(0.2)
    icon_y = top + (height - icon_size) // 2
    page = slide.shapes.add_shape(MSO_SHAPE.FOLDED_CORNER,
                                   int(icon_x), int(icon_y),
                                   int(icon_size), int(icon_size))
    page.fill.solid()
    if pdf_icon:
        page.fill.fore_color.rgb = WHITE
        page.line.color.rgb = NAVY
        page.line.width = Pt(0.75)
    else:
        page.fill.fore_color.rgb = NAVY
        page.line.fill.background()
    page.shadow.inherit = False
    if pdf_icon:
        # "PDF" label inside the white sheet.
        ptf = page.text_frame
        ptf.word_wrap = False
        ptf.margin_left = 0
        ptf.margin_right = 0
        ptf.margin_top = 0
        ptf.margin_bottom = 0
        ptf.vertical_anchor = MSO_ANCHOR.MIDDLE
        pp = ptf.paragraphs[0]
        pp.alignment = PP_ALIGN.CENTER
        pr = pp.add_run()
        pr.text = "PDF"
        pr.font.name = "Calibri"
        pr.font.size = Pt(11)
        pr.font.bold = True
        pr.font.color.rgb = NAVY

    # Label text (italic, navy) inside the card to the right of the icon
    txt_left = int(icon_x + icon_size + Inches(0.2))
    txt_w = int(width - (icon_x + icon_size + Inches(0.4) - left))
    label_box = slide.shapes.add_textbox(txt_left, top,
                                          txt_w, height)
    label_tf = label_box.text_frame
    label_tf.word_wrap = True
    label_tf.margin_left = 0
    label_tf.margin_right = 0
    label_tf.margin_top = 0
    label_tf.margin_bottom = 0
    label_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = label_tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    # First run: "See teaching note:" prefix in caller-supplied color
    # (GOLD by default; NAVY for the slide-33 treatment).
    r1 = p.add_run()
    r1.text = "SEE TEACHING NOTE  →  "
    r1.font.name = "Calibri"
    r1.font.size = Pt(12)
    r1.font.bold = True
    r1.font.color.rgb = label_color
    # Second run: the title of the note, italic navy
    r2 = p.add_run()
    r2.text = text
    r2.font.name = "Calibri"
    r2.font.size = Pt(16)
    r2.font.italic = True
    r2.font.bold = True
    r2.font.color.rgb = NAVY
    return card


def _add_discussion_break(slide, *, top=Inches(6.25), width=Inches(4.8),
                           left=None, text="Discussion Break"):
    """Rounded-parallelogram 'discussion break' badge (bottom-right).

    Custom-geometry shape: top and bottom edges are horizontal; the left
    and right edges slant at 45° in real space (skew = height of the
    shape).  All four corners are slightly rounded.  Gold fill, navy
    bold text, soft drop shadow.

    left=None → pin to the right edge with the default MARGIN. Pass an
    EMU/Inches value to override (used for hand-positioned badges).
    """
    height = Inches(0.72)
    if left is None:
        left = SLIDE_W - MARGIN - width
    left, top, width, height = int(left), int(top), int(width), int(height)
    # Compute skew in path-coordinate units so that left/right sides slant
    # at 45° in REAL space:  horizontal offset of top edge == shape height.
    skew = int(100000 * height / width) if width else 15000
    skew = min(max(skew, 6000), 35000)
    r = 5000          # corner radius in path units — gentle rounding

    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = GOLD
    shp.line.fill.background()
    shp.shadow.inherit = False
    # Strip the default prstGeom – we'll inject custGeom instead.
    spPr = shp._element.spPr
    for old in spPr.findall(qn('a:prstGeom')):
        spPr.remove(old)
    rs = int(r * skew / 100000)
    # IMPORTANT: <a:rect> defines the TEXT bounding rectangle inside the
    # custom geometry.  Set it to the parallelogram's inscribed rectangle
    # (from TL vertex to BR vertex) so PowerPoint won't render text past
    # the slanted edges, regardless of the text frame's own lIns/rIns.
    custgeom_xml = (
        f'<a:custGeom xmlns:a="{A_NS}">'
        f'<a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/>'
        f'<a:rect l="{skew}" t="0" r="{100000-skew}" b="100000"/>'
        f'<a:pathLst><a:path w="100000" h="100000">'
        # Top-left rounded corner (vertex at (skew, 0))
        f'<a:moveTo><a:pt x="{skew+r}" y="0"/></a:moveTo>'
        # Top edge
        f'<a:lnTo><a:pt x="{100000-r}" y="0"/></a:lnTo>'
        # Top-right corner round
        f'<a:cubicBezTo>'
        f'<a:pt x="100000" y="0"/><a:pt x="100000" y="0"/>'
        f'<a:pt x="{100000-rs}" y="{r}"/>'
        f'</a:cubicBezTo>'
        # Right slanted side (down-left)
        f'<a:lnTo><a:pt x="{100000-skew+rs}" y="{100000-r}"/></a:lnTo>'
        # Bottom-right corner round (vertex at (100000-skew, 100000))
        f'<a:cubicBezTo>'
        f'<a:pt x="{100000-skew}" y="100000"/>'
        f'<a:pt x="{100000-skew}" y="100000"/>'
        f'<a:pt x="{100000-skew-r}" y="100000"/>'
        f'</a:cubicBezTo>'
        # Bottom edge
        f'<a:lnTo><a:pt x="{r}" y="100000"/></a:lnTo>'
        # Bottom-left corner round (vertex at (0, 100000))
        f'<a:cubicBezTo>'
        f'<a:pt x="0" y="100000"/><a:pt x="0" y="100000"/>'
        f'<a:pt x="{rs}" y="{100000-r}"/>'
        f'</a:cubicBezTo>'
        # Left slanted side (up-right)
        f'<a:lnTo><a:pt x="{skew-rs}" y="{r}"/></a:lnTo>'
        # Top-left corner round (close the path)
        f'<a:cubicBezTo>'
        f'<a:pt x="{skew}" y="0"/><a:pt x="{skew}" y="0"/>'
        f'<a:pt x="{skew+r}" y="0"/>'
        f'</a:cubicBezTo>'
        f'<a:close/>'
        f'</a:path></a:pathLst>'
        f'</a:custGeom>'
    )
    custgeom = ET.fromstring(custgeom_xml)
    # Insert custGeom right after a:xfrm (schema order)
    xfrm = spPr.find(qn('a:xfrm'))
    if xfrm is not None:
        xfrm.addnext(custgeom)
    else:
        spPr.insert(0, custgeom)

    # Drop shadow (45° down-right, 50% opacity).
    for old in spPr.findall(qn('a:effectLst')):
        spPr.remove(old)
    effectLst = ET.SubElement(spPr, qn('a:effectLst'))
    outerShdw = ET.SubElement(effectLst, qn('a:outerShdw'))
    outerShdw.set('blurRad', '50800')
    outerShdw.set('dist', '38100')
    outerShdw.set('dir', '2700000')
    outerShdw.set('algn', 'tl')
    outerShdw.set('rotWithShape', '0')
    rgb = ET.SubElement(outerShdw, qn('a:srgbClr'))
    rgb.set('val', '000000')
    alpha = ET.SubElement(rgb, qn('a:alpha'))
    alpha.set('val', '50000')

    # Text — render as a SEPARATE textbox overlaid on top of the
    # parallelogram, positioned exactly inside the inscribed rectangle.
    # This decouples text placement from the shape geometry and avoids
    # PowerPoint rendering the run past the slanted edges (which can
    # happen with the in-shape text frame on some PowerPoint versions
    # even with <a:rect> set inside custGeom).
    # In real EMU, the skew equals the shape height (45° slant), so the
    # inscribed rectangle spans (left + height) → (left + width - height).
    skew_emu = height
    ins_left = left + skew_emu
    ins_top = top
    ins_w = width - 2 * skew_emu
    ins_h = height
    txt = slide.shapes.add_textbox(int(ins_left), int(ins_top),
                                     int(ins_w), int(ins_h))
    tf = txt.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(28)         # bumped 20 → 28 on 2026-05-16
    run.font.bold = True
    run.font.color.rgb = NAVY
    return shp


def _add_callout_box(slide, left, top, width, height, text, *,
                      fill=GOLD, text_color=WHITE, size=14, bold=True,
                      rounded=False, shadow=False):
    """Small free-form annotation/callout (e.g., 'plot the slope', 'Revenue
    per car net of material cost').  Used to mark a graph or sub-region.

    ``rounded``/``shadow`` (default off) route to the rounded + drop-shadow
    variant for the deck-wide box treatment."""
    if rounded or shadow:
        return _add_rounded_filled_box(
            slide, left, top, width, height, text,
            fill=fill, text_color=text_color,
            size=size, bold=bold, font="Calibri",
            corner_pct=0.12, shadow=shadow)
    return _add_filled_box(slide, left, top, width, height, text,
                            fill=fill, text_color=text_color,
                            size=size, bold=bold, font="Calibri")


def _add_anchor_burst(slide, left, top, width, height,
                       top_text, bottom_text=None, extra_text=None,
                       *, fill=GOLD, text_color=NAVY,
                       top_size=14, bottom_size=10):
    """12-point star background + a separate text-box overlay.

    The star is purely decorative; the text lives in a normal
    rectangular text box layered on top so the text isn't constrained
    by the star's geometry (no risk of bleeding into the points).
    Reusable on every slide where MB = MC is being invoked, so the same
    visual pattern carries over.
    """
    left, top, width, height = int(left), int(top), int(width), int(height)

    # 1. Decorative star background (no text)
    shp = slide.shapes.add_shape(MSO_SHAPE.STAR_12_POINT, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = NAVY
    shp.line.width = Pt(1.0)
    shp.shadow.inherit = False
    # Soft drop shadow — added 2026-05-15 per user request, so the MB=MC
    # star reads as lifted off the slide like every other content shape.
    _add_drop_shadow(shp)
    # Suppress any auto-inserted text frame contents on the shape itself.
    shp.text_frame.text = ""

    # 2. Overlay text box, sized to the star's inscribed body
    inner_w = int(width * 0.65)
    inner_h = int(height * 0.55)
    inner_x = left + (width - inner_w) // 2
    inner_y = top + (height - inner_h) // 2

    box = slide.shapes.add_textbox(inner_x, inner_y, inner_w, inner_h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = top_text
    r1.font.name = 'Calibri'
    r1.font.size = Pt(top_size)
    r1.font.bold = True
    r1.font.color.rgb = text_color
    if bottom_text:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = bottom_text
        r2.font.name = 'Calibri'
        r2.font.size = Pt(bottom_size)
        r2.font.italic = True
        r2.font.bold = True
        r2.font.color.rgb = text_color
    if extra_text:
        p3 = tf.add_paragraph()
        p3.alignment = PP_ALIGN.CENTER
        r3 = p3.add_run()
        r3.text = extra_text
        r3.font.name = 'Calibri'
        r3.font.size = Pt(bottom_size)
        r3.font.italic = True
        r3.font.bold = True
        r3.font.color.rgb = text_color
    return shp


# --------------------------------------------------------------------------
# OMML (Office Math Markup Language) equation helper – gives formulas a
# proper TeX-style render with italic variables, stacked fractions, real
# subscripts/superscripts.  Uses Cambria Math (the standard PPT math font).
# --------------------------------------------------------------------------

M_NS = 'http://schemas.openxmlformats.org/officeDocument/2006/math'
A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
A14_NS = 'http://schemas.microsoft.com/office/drawing/2010/main'


def _omml_fill(color):
    """Build the inner ``<a:solidFill>`` clause for an OMML run's rPr.

    ``color`` may be an ``RGBColor`` instance (or a 3-tuple of ints).
    Returns an empty string when ``color`` is None so callers can splice
    the result into rPr unconditionally.
    """
    if color is None:
        return ''
    return (
        f'<a:solidFill><a:srgbClr val="'
        f'{color[0]:02X}{color[1]:02X}{color[2]:02X}'
        f'"/></a:solidFill>'
    )


def _omml_run(text, *, color=None):
    """OMML run for an italic variable (default math style).

    Inside an oMath, italic style is the math default for Latin letters;
    we leave m:rPr out entirely so the Cambria Math italic comes through.
    The a:rPr applies drawing-level font sizing/coloring.  Pass ``color``
    to tint the run (e.g., green ΔL / ΔQ in the slide-14 Convention box).
    """
    return (
        f'<m:r xmlns:m="{M_NS}">'
        f'<a:rPr xmlns:a="{A_NS}" lang="en-US" b="0" i="1">'
        f'{_omml_fill(color)}'
        f'<a:latin typeface="Cambria Math"/>'
        f'<a:ea typeface="Cambria Math"/>'
        f'</a:rPr>'
        f'<m:t>{text}</m:t>'
        f'</m:r>'
    )


def _omml_text(text, *, color=None):
    """Upright-style OMML run (for operators, numbers, acronyms).

    Force plain (upright) style via <m:rPr><m:sty m:val="p"/></m:rPr> – this
    is the documented way to disable the math-default italics for the
    enclosed run.  Pass ``color`` to tint the run.
    """
    return (
        f'<m:r xmlns:m="{M_NS}">'
        f'<m:rPr><m:sty m:val="p"/></m:rPr>'
        f'<a:rPr xmlns:a="{A_NS}" lang="en-US" b="0" i="0">'
        f'{_omml_fill(color)}'
        f'<a:latin typeface="Cambria Math"/>'
        f'</a:rPr>'
        f'<m:t xml:space="preserve">{text}</m:t>'
        f'</m:r>'
    )


def _omml_sub(base, sub):
    """OMML subscript: base with subscript expression."""
    return (
        f'<m:sSub xmlns:m="{M_NS}">'
        f'<m:sSubPr><m:ctrlPr>'
        f'<a:rPr xmlns:a="{A_NS}" lang="en-US" i="1">'
        f'<a:latin typeface="Cambria Math"/></a:rPr>'
        f'</m:ctrlPr></m:sSubPr>'
        f'<m:e>{base}</m:e>'
        f'<m:sub>{sub}</m:sub>'
        f'</m:sSub>'
    )


def _omml_frac(num, den):
    """OMML stacked fraction: num / den."""
    return (
        f'<m:f xmlns:m="{M_NS}">'
        f'<m:fPr><m:ctrlPr>'
        f'<a:rPr xmlns:a="{A_NS}" lang="en-US" i="1">'
        f'<a:latin typeface="Cambria Math"/></a:rPr>'
        f'</m:ctrlPr></m:fPr>'
        f'<m:num>{num}</m:num>'
        f'<m:den>{den}</m:den>'
        f'</m:f>'
    )


def _omml_sup(base, sup):
    """OMML superscript: base^sup (e.g. Q²)."""
    return (
        f'<m:sSup xmlns:m="{M_NS}">'
        f'<m:sSupPr><m:ctrlPr>'
        f'<a:rPr xmlns:a="{A_NS}" lang="en-US" i="1">'
        f'<a:latin typeface="Cambria Math"/></a:rPr>'
        f'</m:ctrlPr></m:sSupPr>'
        f'<m:e>{base}</m:e>'
        f'<m:sup>{sup}</m:sup>'
        f'</m:sSup>'
    )


def _add_dashed_gridlines(axis_el):
    """Dashed light-grey major gridlines on an axis (Cobb-Douglas style)."""
    for old in axis_el.findall(qn('c:majorGridlines')):
        axis_el.remove(old)
    gl = ET.Element(qn('c:majorGridlines'))
    sp = ET.SubElement(gl, qn('c:spPr'))
    ln = ET.SubElement(sp, qn('a:ln'))
    ln.set('w', '9525')
    ln.set('cap', 'flat'); ln.set('cmpd', 'sng'); ln.set('algn', 'ctr')
    fill = ET.SubElement(ln, qn('a:solidFill'))
    clr = ET.SubElement(fill, qn('a:srgbClr')); clr.set('val', 'C8CDD3')
    dash = ET.SubElement(ln, qn('a:prstDash')); dash.set('val', 'dash')
    axpos = axis_el.find(qn('c:axPos'))
    if axpos is not None:
        axpos.addnext(gl)


def _align_x_labels_with_ticks(value_axis):
    """Set <c:crossBetween val="midCat"/> on a value axis so the category
    labels and tick marks align (default for line charts in OOXML is
    "between", which leaves labels visually between adjacent ticks).
    """
    val_el = value_axis._element
    for old in val_el.findall(qn('c:crossBetween')):
        val_el.remove(old)
    cb = ET.Element(qn('c:crossBetween'))
    cb.set('val', 'midCat')
    # Schema position: c:crossBetween follows c:crosses(At) and precedes
    # c:majorUnit.  Insert before c:majorUnit if present; else append.
    mu_el = val_el.find(qn('c:majorUnit'))
    if mu_el is not None:
        mu_el.addprevious(cb)
    else:
        val_el.append(cb)


def _make_simple_line_chart(slide, x, y, w, h, categories, values, *,
                              line_color, x_title, y_title,
                              y_min=0, y_max=None, y_unit=None,
                              marker='circle'):
    """Single-series line+markers chart with dashed light-grey gridlines.

    Same visual conventions as slide 11 (Calibri navy labels, dashed C8CDD3
    major gridlines, marker size 7) but no legend/title.
    """
    cd = CategoryChartData()
    cd.categories = list(categories)
    cd.add_series("Y", values)
    # Drop-shadow rectangle behind the chart (graphicFrames can't host shadow).
    _add_graphicframe_shadow(slide, x, y, w, h)
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE,
        int(x), int(y), int(w), int(h), cd,
    )
    chart = chart_shape.chart
    chart.has_title = False
    chart.has_legend = False

    clr_hex = f'{line_color[0]:02X}{line_color[1]:02X}{line_color[2]:02X}'

    for series in chart.series:
        line = series.format.line
        line.color.rgb = line_color
        line.width = Pt(2.5)
        ser_xml = series._element
        for old in ser_xml.findall(qn('c:marker')):
            ser_xml.remove(old)
        m = ET.SubElement(ser_xml, qn('c:marker'))
        sym = ET.SubElement(m, qn('c:symbol')); sym.set('val', marker)
        sz_el = ET.SubElement(m, qn('c:size')); sz_el.set('val', '7')
        sp = ET.SubElement(m, qn('c:spPr'))
        fl = ET.SubElement(sp, qn('a:solidFill'))
        rg = ET.SubElement(fl, qn('a:srgbClr')); rg.set('val', clr_hex)
        ln = ET.SubElement(sp, qn('a:ln'))
        lf = ET.SubElement(ln, qn('a:solidFill'))
        lr = ET.SubElement(lf, qn('a:srgbClr')); lr.set('val', clr_hex)
        # disable smoothing (straight segments between points)
        for sm in ser_xml.findall(qn('c:smooth')):
            ser_xml.remove(sm)
        smooth = ET.SubElement(ser_xml, qn('c:smooth'))
        smooth.set('val', '0')

    # Axes – axis titles in BOLD ITALIC navy (per course CLAUDE.md);
    # tick labels in regular Calibri navy.
    cat = chart.category_axis
    cat.tick_labels.font.name = "Calibri"
    cat.tick_labels.font.size = Pt(10)
    cat.tick_labels.font.color.rgb = NAVY
    cat.has_title = True
    cat.axis_title.text_frame.text = x_title
    ar = cat.axis_title.text_frame.paragraphs[0].runs[0]
    ar.font.name = "Calibri"; ar.font.size = Pt(12)
    ar.font.bold = True; ar.font.italic = True
    ar.font.color.rgb = NAVY

    val = chart.value_axis
    val.tick_labels.font.name = "Calibri"
    val.tick_labels.font.size = Pt(10)
    val.tick_labels.font.color.rgb = NAVY
    val.has_title = True
    val.axis_title.text_frame.text = y_title
    ar = val.axis_title.text_frame.paragraphs[0].runs[0]
    ar.font.name = "Calibri"; ar.font.size = Pt(12)
    ar.font.bold = True; ar.font.italic = True
    ar.font.color.rgb = NAVY
    if y_max is not None:
        val.minimum_scale = y_min
        val.maximum_scale = y_max
    if y_unit is not None:
        val.major_unit = y_unit

    # Align X-axis category labels with tick marks (default OOXML places
    # them in the gaps between ticks).
    _align_x_labels_with_ticks(val)

    _add_dashed_gridlines(cat._element)
    _add_dashed_gridlines(val._element)
    return chart_shape


def _make_multi_line_chart(slide, x, y, w, h, categories, series, *,
                             x_title, y_title,
                             y_min=0, y_max=None, y_unit=None,
                             legend=True, legend_pos=('0.08', '0.10', '0.20', '0.20')):
    """Multi-series line+markers chart with the deck's standard styling.

    series: list of (name, values, color: RGBColor, marker: str) tuples.
    legend_pos: (x, y, w, h) in chart-fraction units (str) – top-left default.
    """
    _add_graphicframe_shadow(slide, x, y, w, h)
    cd = CategoryChartData()
    cd.categories = list(categories)
    for name, values, _color, _marker in series:
        cd.add_series(name, values)
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE,
        int(x), int(y), int(w), int(h), cd,
    )
    chart = chart_shape.chart
    chart.has_title = False
    chart.has_legend = legend

    # Per-series styling: line color + marker
    for idx, ser in enumerate(chart.series):
        name, values, color, marker = series[idx]
        line = ser.format.line
        line.color.rgb = color
        line.width = Pt(2.5)
        ser_xml = ser._element
        clr_hex = f'{color[0]:02X}{color[1]:02X}{color[2]:02X}'
        for old in ser_xml.findall(qn('c:marker')):
            ser_xml.remove(old)
        m = ET.SubElement(ser_xml, qn('c:marker'))
        sym = ET.SubElement(m, qn('c:symbol')); sym.set('val', marker)
        sz_el = ET.SubElement(m, qn('c:size')); sz_el.set('val', '7')
        sp = ET.SubElement(m, qn('c:spPr'))
        fl = ET.SubElement(sp, qn('a:solidFill'))
        rg = ET.SubElement(fl, qn('a:srgbClr')); rg.set('val', clr_hex)
        ln = ET.SubElement(sp, qn('a:ln'))
        lf = ET.SubElement(ln, qn('a:solidFill'))
        lr = ET.SubElement(lf, qn('a:srgbClr')); lr.set('val', clr_hex)
        for sm in ser_xml.findall(qn('c:smooth')):
            ser_xml.remove(sm)
        smooth = ET.SubElement(ser_xml, qn('c:smooth'))
        smooth.set('val', '0')

    # Axes – bold italic navy titles per course style
    cat = chart.category_axis
    cat.tick_labels.font.name = "Calibri"
    cat.tick_labels.font.size = Pt(10)
    cat.tick_labels.font.color.rgb = NAVY
    cat.has_title = True
    cat.axis_title.text_frame.text = x_title
    ar = cat.axis_title.text_frame.paragraphs[0].runs[0]
    ar.font.name = "Calibri"; ar.font.size = Pt(12)
    ar.font.bold = True; ar.font.italic = True; ar.font.color.rgb = NAVY

    val = chart.value_axis
    val.tick_labels.font.name = "Calibri"
    val.tick_labels.font.size = Pt(10)
    val.tick_labels.font.color.rgb = NAVY
    val.has_title = True
    val.axis_title.text_frame.text = y_title
    ar = val.axis_title.text_frame.paragraphs[0].runs[0]
    ar.font.name = "Calibri"; ar.font.size = Pt(12)
    ar.font.bold = True; ar.font.italic = True; ar.font.color.rgb = NAVY
    if y_max is not None:
        val.minimum_scale = y_min
        val.maximum_scale = y_max
    if y_unit is not None:
        val.major_unit = y_unit

    # Align X-axis category labels with tick marks.
    _align_x_labels_with_ticks(val)

    _add_dashed_gridlines(cat._element)
    _add_dashed_gridlines(val._element)

    # Legend top-left inside plot with white fill
    if legend:
        leg_el = chart.legend._element
        chart.legend.font.name = "Calibri"
        chart.legend.font.size = Pt(11)
        chart.legend.font.color.rgb = NAVY
        chart.legend.include_in_layout = False
        # Strip default legendPos / layout, replace
        for old in leg_el.findall(qn('c:layout')):
            leg_el.remove(old)
        for old in leg_el.findall(qn('c:legendPos')):
            leg_el.remove(old)
        pos = ET.SubElement(leg_el, qn('c:legendPos')); pos.set('val', 'tr')
        leg_el.remove(pos); leg_el.insert(0, pos)
        # manualLayout positions in chart-fraction units
        layout = ET.Element(qn('c:layout'))
        ml = ET.SubElement(layout, qn('c:manualLayout'))
        ET.SubElement(ml, qn('c:xMode')).set('val', 'edge')
        ET.SubElement(ml, qn('c:yMode')).set('val', 'edge')
        ET.SubElement(ml, qn('c:x')).set('val', legend_pos[0])
        ET.SubElement(ml, qn('c:y')).set('val', legend_pos[1])
        ET.SubElement(ml, qn('c:w')).set('val', legend_pos[2])
        ET.SubElement(ml, qn('c:h')).set('val', legend_pos[3])
        pos.addnext(layout)
        # White fill behind legend
        for old in leg_el.findall(qn('c:spPr')):
            leg_el.remove(old)
        leg_spPr = ET.Element(qn('c:spPr'))
        sf = ET.SubElement(leg_spPr, qn('a:solidFill'))
        clr = ET.SubElement(sf, qn('a:srgbClr')); clr.set('val', 'FFFFFF')
        ln = ET.SubElement(leg_spPr, qn('a:ln')); ln.set('w', '6350')
        lf = ET.SubElement(ln, qn('a:solidFill'))
        lc = ET.SubElement(lf, qn('a:srgbClr')); lc.set('val', '0B2B4E')
        layout.addnext(leg_spPr)
    return chart_shape


def _make_xy_line_chart(slide, x, y, w, h, *, series, x_title, y_title,
                          x_min=0, x_max=None, x_unit=None,
                          y_min=0, y_max=None, y_unit=None,
                          legend=False, legend_pos=None, smooth=False):
    """XY-scatter line+markers chart with the deck's standard styling.

    Use when you need data points to sit at arbitrary X-positions (e.g.,
    plotting MPL at the midpoint of each L-interval) while keeping tick
    marks at standard L-values.  Both axes are numeric value axes.

    series: list of ``(name, [(x, y), ...], color: RGBColor, marker: str)``.
    smooth: True for XY_SCATTER_SMOOTH (cubic spline through points),
            False for XY_SCATTER_LINES (straight segments).
    """
    _add_graphicframe_shadow(slide, x, y, w, h)
    cd = XyChartData()
    for name, points, _color, _marker in series:
        s = cd.add_series(name)
        for px, py in points:
            s.add_data_point(px, py)
    chart_type = (XL_CHART_TYPE.XY_SCATTER_SMOOTH if smooth
                  else XL_CHART_TYPE.XY_SCATTER_LINES)
    chart_shape = slide.shapes.add_chart(
        chart_type,
        int(x), int(y), int(w), int(h), cd,
    )
    chart = chart_shape.chart
    chart.has_title = False
    chart.has_legend = legend

    # Per-series styling
    for idx, ser in enumerate(chart.series):
        name, points, color, marker = series[idx]
        line = ser.format.line
        line.color.rgb = color
        line.width = Pt(2.5)
        ser_xml = ser._element
        clr_hex = f'{color[0]:02X}{color[1]:02X}{color[2]:02X}'
        for old in ser_xml.findall(qn('c:marker')):
            ser_xml.remove(old)
        m = ET.SubElement(ser_xml, qn('c:marker'))
        sym = ET.SubElement(m, qn('c:symbol')); sym.set('val', marker)
        sz_el = ET.SubElement(m, qn('c:size')); sz_el.set('val', '7')
        sp = ET.SubElement(m, qn('c:spPr'))
        fl = ET.SubElement(sp, qn('a:solidFill'))
        rg = ET.SubElement(fl, qn('a:srgbClr')); rg.set('val', clr_hex)
        ln = ET.SubElement(sp, qn('a:ln'))
        lf = ET.SubElement(ln, qn('a:solidFill'))
        lr = ET.SubElement(lf, qn('a:srgbClr')); lr.set('val', clr_hex)
        for sm in ser_xml.findall(qn('c:smooth')):
            ser_xml.remove(sm)
        smooth = ET.SubElement(ser_xml, qn('c:smooth'))
        smooth.set('val', '0')

    # Both axes are value axes in XY scatter; python-pptx returns the
    # X axis through .category_axis (wrapped as a ValueAxis since
    # catAx_lst is empty) and the Y axis through .value_axis.
    x_ax = chart.category_axis
    x_ax.tick_labels.font.name = "Calibri"
    x_ax.tick_labels.font.size = Pt(10)
    x_ax.tick_labels.font.color.rgb = NAVY
    if x_max is not None:
        x_ax.minimum_scale = x_min
        x_ax.maximum_scale = x_max
    if x_unit is not None:
        x_ax.major_unit = x_unit
    x_ax.has_title = True
    x_ax.axis_title.text_frame.text = x_title
    ar = x_ax.axis_title.text_frame.paragraphs[0].runs[0]
    ar.font.name = "Calibri"; ar.font.size = Pt(12)
    ar.font.bold = True; ar.font.italic = True; ar.font.color.rgb = NAVY

    y_ax = chart.value_axis
    y_ax.tick_labels.font.name = "Calibri"
    y_ax.tick_labels.font.size = Pt(10)
    y_ax.tick_labels.font.color.rgb = NAVY
    if y_max is not None:
        y_ax.minimum_scale = y_min
        y_ax.maximum_scale = y_max
    if y_unit is not None:
        y_ax.major_unit = y_unit
    y_ax.has_title = True
    y_ax.axis_title.text_frame.text = y_title
    ar = y_ax.axis_title.text_frame.paragraphs[0].runs[0]
    ar.font.name = "Calibri"; ar.font.size = Pt(12)
    ar.font.bold = True; ar.font.italic = True; ar.font.color.rgb = NAVY

    _add_dashed_gridlines(x_ax._element)
    _add_dashed_gridlines(y_ax._element)

    # Legend
    if legend and legend_pos is not None:
        leg_el = chart.legend._element
        chart.legend.font.name = "Calibri"
        chart.legend.font.size = Pt(11)
        chart.legend.font.color.rgb = NAVY
        chart.legend.include_in_layout = False
        for old in leg_el.findall(qn('c:layout')):
            leg_el.remove(old)
        for old in leg_el.findall(qn('c:legendPos')):
            leg_el.remove(old)
        pos = ET.SubElement(leg_el, qn('c:legendPos')); pos.set('val', 'tr')
        leg_el.remove(pos); leg_el.insert(0, pos)
        layout = ET.Element(qn('c:layout'))
        ml = ET.SubElement(layout, qn('c:manualLayout'))
        ET.SubElement(ml, qn('c:xMode')).set('val', 'edge')
        ET.SubElement(ml, qn('c:yMode')).set('val', 'edge')
        ET.SubElement(ml, qn('c:x')).set('val', legend_pos[0])
        ET.SubElement(ml, qn('c:y')).set('val', legend_pos[1])
        ET.SubElement(ml, qn('c:w')).set('val', legend_pos[2])
        ET.SubElement(ml, qn('c:h')).set('val', legend_pos[3])
        pos.addnext(layout)
        for old in leg_el.findall(qn('c:spPr')):
            leg_el.remove(old)
        leg_spPr = ET.Element(qn('c:spPr'))
        sf = ET.SubElement(leg_spPr, qn('a:solidFill'))
        clr = ET.SubElement(sf, qn('a:srgbClr')); clr.set('val', 'FFFFFF')
        ln = ET.SubElement(leg_spPr, qn('a:ln')); ln.set('w', '6350')
        lf = ET.SubElement(ln, qn('a:solidFill'))
        lc = ET.SubElement(lf, qn('a:srgbClr')); lc.set('val', '0B2B4E')
        layout.addnext(leg_spPr)

    return chart_shape


def _omml_acc_overline(symbol):
    """Inline OMML accent (overline / bar) on a math symbol.

    symbol: e.g. 'K' or 'L' – the variable to wear the bar.
    Returns an OMML fragment intended to be embedded inside an <m:oMath>.
    """
    return (
        '<m:acc>'
          '<m:accPr><m:chr m:val="̅"/></m:accPr>'
          '<m:e>'
            '<m:r>'
              '<a:rPr lang="en-US" b="0" i="1">'
                '<a:latin typeface="Cambria Math"/>'
              '</a:rPr>'
              f'<m:t>{symbol}</m:t>'
            '</m:r>'
          '</m:e>'
        '</m:acc>'
    )


def _add_mixed_textbox(slide, left, top, width, height, segments, *,
                        align=PP_ALIGN.LEFT, default_color=NAVY,
                        default_size=24,
                        margin_left=None, margin_right=None,
                        margin_top=None, margin_bottom=None):
    """Build a textbox whose paragraphs mix plain text runs and inline OMML.

    segments: list of (kind, content, opts) tuples, with kind ∈ {"text",
    "omml", "break"}.  "break" inserts a new paragraph.  Opts may set
    `size`, `bold`, `italic`, `color`, `font` per run.
    """
    box = slide.shapes.add_textbox(int(left), int(top),
                                     int(width), int(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05) if margin_left is None else margin_left
    tf.margin_right = Inches(0.05) if margin_right is None else margin_right
    tf.margin_top = Inches(0.0) if margin_top is None else margin_top
    tf.margin_bottom = Inches(0.0) if margin_bottom is None else margin_bottom

    align_attr = ''
    if align == PP_ALIGN.CENTER: align_attr = ' algn="ctr"'
    elif align == PP_ALIGN.RIGHT: align_attr = ' algn="r"'

    def _start_para():
        return [f'<a:p xmlns:a="{A_NS}" xmlns:m="{M_NS}" xmlns:a14="{A14_NS}">',
                f'<a:pPr{align_attr}/>' if align_attr else '']

    paragraphs = [_start_para()]
    for kind, content, opts in segments:
        if kind == 'break':
            paragraphs[-1].append('<a:endParaRPr lang="en-US"/></a:p>')
            paragraphs.append(_start_para())
            continue
        size_pt = int(opts.get('size', default_size) * 100)
        color = opts.get('color', default_color)
        clr_hex = f'{color[0]:02X}{color[1]:02X}{color[2]:02X}'
        bold_attr = ' b="1"' if opts.get('bold') else ''
        italic_attr = ' i="1"' if opts.get('italic') else ''
        font = opts.get('font', 'Calibri')
        if kind == 'text':
            paragraphs[-1].append(
                f'<a:r><a:rPr lang="en-US" sz="{size_pt}"{bold_attr}{italic_attr}>'
                f'<a:solidFill><a:srgbClr val="{clr_hex}"/></a:solidFill>'
                f'<a:latin typeface="{font}"/>'
                f'</a:rPr><a:t>{content}</a:t></a:r>'
            )
        elif kind == 'omml':
            paragraphs[-1].append(
                f'<a14:m><m:oMath>{content}</m:oMath></a14:m>'
            )
    paragraphs[-1].append('<a:endParaRPr lang="en-US"/></a:p>')
    full_xml = ''.join(''.join(p) for p in paragraphs)

    txBody = tf._txBody
    for old in list(txBody.findall(qn('a:p'))):
        txBody.remove(old)
    # We have to parse each <a:p> separately so the namespaces resolve
    for p_str in full_xml.split('</a:p>'):
        if not p_str.strip(): continue
        p_xml = p_str + '</a:p>'
        new_p = ET.fromstring(p_xml)
        txBody.append(new_p)
        # Apply size+color to any OMML m:r elements inside this paragraph.
        # Color is set to ``default_color`` ONLY when no per-run solidFill
        # is already present — this lets callers tint individual OMML runs
        # via the optional ``color=`` argument on ``_omml_run`` / ``_omml_text``
        # (e.g., the green ΔL / ΔQ in the slide-14 Convention box) without
        # being silently overridden here.
        clr_hex = f'{default_color[0]:02X}{default_color[1]:02X}{default_color[2]:02X}'
        for r in new_p.iter(qn('m:r')):
            arPr = r.find(qn('a:rPr'))
            if arPr is None:
                arPr = ET.Element(qn('a:rPr'))
                arPr.set('lang', 'en-US')
                r.insert(0, arPr)
            arPr.set('sz', str(int(default_size * 100)))
            if arPr.find(qn('a:solidFill')) is None:
                # fill BEFORE a:latin (schema order) or the color is ignored
                sf = arPr.makeelement(qn('a:solidFill'), {})
                srgb = ET.SubElement(sf, qn('a:srgbClr'))
                srgb.set('val', clr_hex)
                arPr.insert(0, sf)
    return box


def _add_math_equation(slide, left, top, width, height, omml_content, *,
                       size_pt=32, color=NAVY, fill=None, line=None,
                       rounded=False, shadow=False, corner_pct=25000):
    """Place an OMML equation in a textbox on the slide.

    omml_content: a string built from _omml_* helpers (without the outer
    <m:oMathPara> wrapper).

    ``rounded=True`` gives the fill box rounded corners (``corner_pct`` in
    OOXML adj units, 25000 ≈ 25%); ``shadow=True`` adds a soft drop
    shadow.  Both default off.  IMPORTANT: the roundRect prstGeom must be
    inserted right after <a:xfrm> and BEFORE the fill, or PowerPoint
    silently drops the shape (see slide-54 fix).
    """
    left, top, width, height = int(left), int(top), int(width), int(height)
    box = slide.shapes.add_textbox(left, top, width, height)

    if fill is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
    if line is not None:
        box.line.color.rgb = line
        box.line.width = Pt(0.75)

    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Replace the default <a:p> with one that hosts a math zone (<a14:m>)
    # containing the oMathPara.  The <a14:m> wrapper is REQUIRED by
    # PowerPoint to recognise OMML inside a textbox – without it PPT just
    # shows empty boxes.
    txBody = tf._txBody
    for p in list(txBody.findall(qn('a:p'))):
        txBody.remove(p)

    sz = int(size_pt * 100)
    clr_hex = '{:02X}{:02X}{:02X}'.format(color[0], color[1], color[2])

    # Use unique namespace prefixes; lxml will resolve them when parsing.
    p_xml = (
        f'<a:p xmlns:a="{A_NS}" xmlns:m="{M_NS}" xmlns:a14="{A14_NS}">'
        f'<a:pPr algn="ctr">'
        f'<a:defRPr sz="{sz}" b="0" i="1">'
        f'<a:solidFill><a:srgbClr val="{clr_hex}"/></a:solidFill>'
        f'<a:latin typeface="Cambria Math"/>'
        f'</a:defRPr>'
        f'</a:pPr>'
        f'<a14:m>'
        f'<m:oMathPara>'
        f'<m:oMathParaPr><m:jc m:val="centerGroup"/></m:oMathParaPr>'
        f'<m:oMath>{omml_content}</m:oMath>'
        f'</m:oMathPara>'
        f'</a14:m>'
        f'<a:endParaRPr lang="en-US" sz="{sz}"/>'
        f'</a:p>'
    )

    new_p = ET.fromstring(p_xml)
    txBody.append(new_p)

    # Set size and color on every OMML run's <a:rPr> so the text renders at
    # the right size.  (The defRPr above is the fallback.)
    for r in new_p.iter(qn('m:r')):
        arPr = r.find(qn('a:rPr'))
        if arPr is None:
            arPr = ET.Element(qn('a:rPr'))
            r.insert(0, arPr)
        arPr.set('sz', str(sz))
        if arPr.get('lang') is None:
            arPr.set('lang', 'en-US')
        # respect per-run colors set via _omml_run/_omml_text(color=...);
        # only runs WITHOUT a fill get the equation's default color
        # (was clobbering firm-color runs — Nico 2026-08-07)
        if not arPr.findall(qn('a:solidFill')):
            sf = ET.Element(qn('a:solidFill'))
            srgb = ET.SubElement(sf, qn('a:srgbClr'))
            srgb.set('val', clr_hex)
            arPr.insert(0, sf)

    # Optional rounded corners + drop shadow on the fill box.  prstGeom
    # MUST sit right after <a:xfrm> and before <a:solidFill>, else PPT
    # silently refuses to render the shape (slide-54 lesson).
    if rounded or shadow:
        spPr = box._element.find(qn('p:spPr'))
        if spPr is not None:
            if rounded:
                for old in spPr.findall(qn('a:prstGeom')):
                    spPr.remove(old)
                prstGeom = ET.Element(qn('a:prstGeom'))
                prstGeom.set('prst', 'roundRect')
                avLst = ET.SubElement(prstGeom, qn('a:avLst'))
                gd = ET.SubElement(avLst, qn('a:gd'))
                gd.set('name', 'adj'); gd.set('fmla', f'val {int(corner_pct)}')
                xfrm = spPr.find(qn('a:xfrm'))
                if xfrm is not None:
                    xfrm.addnext(prstGeom)
                else:
                    spPr.insert(0, prstGeom)
            if shadow:
                _add_drop_shadow(box)
    return box




# ==========================================================================
#  MODULE 7 — OLIGOPOLY AND GAME THEORY  (everything below is M7-specific)
#
#  Above this line: the reusable helper layer, copied VERBATIM from
#  Module 3/_build_Module3.py (palette, chrome, boxes, badges, bullets,
#  OMML engine, chart/table primitives).  Do not edit the helpers here —
#  improvements belong upstream in a shared module someday.
#
#  Build phases (course CLAUDE.md):
#    Phase 1 (THIS SCRIPT): all script-buildable slides. Poll/video slides
#      are stubs holding their position; hidden slides carry show="0".
#    Phase 2: freeze this script once all buildable slides are approved.
#    Phase 3: OOXML surgery — splice the 5 PollEv + 3 video slides from
#      "Module 7.pptx" (slides 32, 52, 69, 72, 77 / 14, 74, 76+78 with
#      their NULL-external video rels), grouping + animation passes.
# ==========================================================================

import uuid

CREAM = RGBColor(0xFD, 0xF6, 0xE6)

TAG_RECAP   = "Module 7 · Recap"
TAG_ROADMAP = "Module 7 · Course Roadmap"
TAG_OUTLINE = "Module 7 · Outline"
TAG_OLIGO   = "Module 7 · Part 1 · Oligopoly"
TAG_COLLU   = "Module 7 · Part 1 · Collusion and Cartels"
TAG_COURNOT = "Module 7 · Part 1 · Cournot Competition"
TAG_BERTRAND= "Module 7 · Part 1 · Bertrand Competition"
TAG_DIFF    = "Module 7 · Part 1 · Differentiated Goods"
TAG_GT      = "Module 7 · Part 2 · Key Concepts"
TAG_GAMES   = "Module 7 · Part 2 · Classic Games"
TAG_COMMIT  = "Module 7 · Part 2 · Commitment"


# --------------------------------------------------------------------------
# M7 override: footer page number is a LIVE slide-number field
# (Teaching CLAUDE.md: "Footer page numbers are LIVE slide-number fields").
# Defined AFTER the helpers so every helper call resolves to this version.
# --------------------------------------------------------------------------

def _add_slidenum_field(slide, left, top, width, height, page_num, *,
                        size=12, color=GRAY):
    tb = slide.shapes.add_textbox(int(left), int(top), int(width), int(height))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    p_el = p._p
    guid = str(uuid.uuid5(uuid.NAMESPACE_DNS, f"m7-slidenum-{page_num}")).upper()
    fld = p_el.makeelement(qn('a:fld'), {'id': '{%s}' % guid, 'type': 'slidenum'})
    rPr = p_el.makeelement(qn('a:rPr'),
                           {'lang': 'en-US', 'sz': str(size * 100), 'dirty': '0'})
    fill = p_el.makeelement(qn('a:solidFill'), {})
    clr = p_el.makeelement(qn('a:srgbClr'), {'val': str(color)})
    fill.append(clr)
    rPr.append(fill)
    latin = p_el.makeelement(qn('a:latin'), {'typeface': 'Calibri'})
    rPr.append(latin)
    t = p_el.makeelement(qn('a:t'), {})
    t.text = str(page_num)
    fld.append(rPr)
    fld.append(t)
    p_el.append(fld)
    return tb


def _draw_footer(slide, footer_text, page_num):  # noqa: F811 — M7 override
    _add_rect(slide, 0, Inches(7.15), SLIDE_W, Inches(0.02), RULE)
    _add_rect(slide, MARGIN, Inches(7.135), GOLD_W, Inches(0.05), GOLD)
    _add_text(slide, MARGIN, Inches(7.20), Inches(11), Inches(0.32),
              footer_text, size=12, color=GRAY)
    _add_slidenum_field(slide, Inches(12.55), Inches(7.20), Inches(0.55),
                        Inches(0.32), page_num)


# --------------------------------------------------------------------------
# M7 image loader — source images live in _source_images/ under their
# original media names (imageNN.ext), extracted 2026-07-29.
# --------------------------------------------------------------------------

def _add_media_image(slide, fname, *, left, top, width=None, height=None,
                     rounded=True, shadow=True, corner_pct=8):
    """Place a source-deck image by media filename. Logos/screenshots:
    rounded=False, shadow=False (flat exception)."""
    path = SRC_IMG_DIR / fname
    kwargs = {"left": int(left), "top": int(top)}
    if width is not None:
        kwargs["width"] = int(width)
    if height is not None:
        kwargs["height"] = int(height)
    pic = slide.shapes.add_picture(str(path), **kwargs)
    if rounded:
        _apply_picture_style(pic, corner_pct=corner_pct)
    elif shadow:
        _add_drop_shadow(pic)
    return pic


# --------------------------------------------------------------------------
# Native styled table (navy header, white/cream body, thin borders) on a
# shadowed white backing card.  NOTE: backing + table must be GROUPED in
# phase 3 (grpSp surgery) so the shade travels with the table.
# --------------------------------------------------------------------------

def _set_cell_borders(cell, *, color=RULE, weight_pt=0.75):
    tcPr = cell._tc.get_or_add_tcPr()
    for tag in ('a:lnL', 'a:lnR', 'a:lnT', 'a:lnB'):
        for old in tcPr.findall(qn(tag)):
            tcPr.remove(old)
    # schema order: lnL, lnR, lnT, lnB come first in tcPr
    for tag in reversed(('a:lnL', 'a:lnR', 'a:lnT', 'a:lnB')):
        ln = tcPr.makeelement(qn(tag), {'w': str(int(weight_pt * 12700)),
                                        'cap': 'flat'})
        fill = ln.makeelement(qn('a:solidFill'), {})
        c = fill.makeelement(qn('a:srgbClr'), {'val': str(color)})
        fill.append(c)
        ln.append(fill)
        tcPr.insert(0, ln)


def _add_styled_table(slide, left, top, width, height, rows_data, *,
                      col_widths=None, row_heights=None, font_size=18,
                      header_size=None, backing_pad=Inches(0.15),
                      first_col_bold=True, first_col_align_left=True,
                      cell_fills=None, cell_text_colors=None):
    """rows_data: list of rows (row 0 = navy header). cell_fills /
    cell_text_colors: optional {(r, c): RGBColor} overrides."""
    header_size = header_size or font_size
    left, top, width, height = int(left), int(top), int(width), int(height)
    _add_graphicframe_shadow(slide, left - int(backing_pad),
                             top - int(backing_pad),
                             width + 2 * int(backing_pad),
                             height + 2 * int(backing_pad))
    n_rows, n_cols = len(rows_data), len(rows_data[0])
    gf = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = gf.table
    tblPr = tbl._tbl.find(qn('a:tblPr'))
    if tblPr is not None:                      # strip default banding style
        tblPr.set('firstRow', '0')
        tblPr.set('bandRow', '0')
        for child in list(tblPr):
            tblPr.remove(child)
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = int(w)
    if row_heights:
        for i, h in enumerate(row_heights):
            tbl.rows[i].height = int(h)
    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.margin_left = Inches(0.10)
            cell.margin_right = Inches(0.10)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            hdr = (r == 0)
            if cell_fills and (r, c) in cell_fills:
                cell.fill.solid()
                cell.fill.fore_color.rgb = cell_fills[(r, c)]
            elif hdr:
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else CREAM
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            if first_col_align_left and c == 0 and not hdr:
                p.alignment = PP_ALIGN.LEFT
            else:
                p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(val)
            run.font.name = "Calibri"
            run.font.size = Pt(header_size if hdr else font_size)
            run.font.bold = hdr or (first_col_bold and c == 0)
            if cell_text_colors and (r, c) in cell_text_colors:
                run.font.color.rgb = cell_text_colors[(r, c)]
            else:
                run.font.color.rgb = WHITE if hdr else NAVY
            _set_cell_borders(cell)
    return gf


# --------------------------------------------------------------------------
# M7 title slide, outline/divider, and stub
# --------------------------------------------------------------------------

def _inject_handoff_group(slide, fname, id_base=9500):
    """Append a hand-authored <p:grpSp> (saved verbatim from a canonical
    slide into _handoff_*.xml) into this slide's shape tree. Shape ids are
    re-based to avoid collisions. Used for Nico's scaled group cards
    (group scaling also scales text, which a rebuild can't replicate)."""
    xml = (OUT_DIR / fname).read_text(encoding='utf-8')
    el = ET.fromstring(xml)
    for i, nv in enumerate(el.iter(qn('p:cNvPr'))):
        nv.set('id', str(id_base + i))
    slide.shapes._spTree.append(el)
    return el


def slide_01_title(prs):
    slide = _blank_slide(prs)
    # Text block shifted up/left and Dogopoly cartoon added on the right
    # — hand-edited by Nico in PowerPoint, ported 2026-08-03 (exact EMUs
    # from the canonical deck; blob = image13.jpg from the source deck).
    tx_l, tx_w = -168626, 11687697
    _add_text(slide, tx_l, 1152144, tx_w, Inches(1.3),
              "Oligopoly and Game Theory",
              size=60, bold=True, color=NAVY, font="Calibri",
              align=PP_ALIGN.CENTER)
    _add_text(slide, tx_l, 2203704, tx_w, Inches(0.75),
              "Module 7",
              size=40, bold=True, color=GOLD, font="Calibri",
              align=PP_ALIGN.CENTER)
    _add_rect(slide, 3846422, 3163824, 3657600, 54864, GOLD)
    _add_text(slide, tx_l, 3529584, tx_w, Inches(0.55),
              "Management 405", size=26, bold=True, color=GRAY,
              font="Calibri", align=PP_ALIGN.CENTER)
    _add_text(slide, tx_l, 4169664, tx_w, Inches(0.5),
              "Prof. Nico Voigtländer  ·  UCLA Anderson",
              size=22, color=GRAY, font="Calibri", align=PP_ALIGN.CENTER)
    _add_rect(slide, 0, Inches(7.15), SLIDE_W, Inches(0.02), RULE)
    _add_rect(slide, MARGIN, Inches(7.135), GOLD_W, Inches(0.05), GOLD)
    _add_media_image(slide, "image13.jpg", left=8248214, top=2189988,
                     width=3871436, height=4267181,
                     rounded=False, shadow=False)
    return slide


# M7 flat outline items for the numbered-circle agenda format (Teaching
# CLAUDE.md "Module-Outline / Agenda Slides"), adopted from Module 2 on
# 2026-08-20 (Nico approved the flattening + descriptions). The old
# part/sub two-level agenda was replaced wholesale.
M7_OUTLINE = [
    ("Collusion and Cartels",
     "Why firms are tempted to collude, and why cartels break down"),
    ("Cournot Competition",
     "Competing by choosing quantities: capacity decisions in oligopoly"),
    ("Bertrand Competition",
     "Competing by setting prices, and why it can compete profits away"),
    ("Oligopoly with Differentiated Goods",
     "How differentiation softens price competition"),
    ("Strategic Thinking: Key Concepts",
     "Best responses, dominant strategies, and Nash equilibrium"),
    ("Prisoner’s Dilemma, Game of Chicken…",
     "Classic games and what they teach about business strategy"),
]


def make_module_outline(prs, page_num, *, section_tag=TAG_OUTLINE,
                        part_idx=None, sub_idx=None, descriptions=False,
                        title="Outline of Module 7"):
    """Module outline in the numbered-circle format (gold 0.58" circle at
    x=1.15, 25 pt bold navy number + title, 22 pt gray description).
    Every item RESERVES the description row (uniform pitch, so item
    positions are identical on every agenda slide); the description text
    shows only for the current topic, or for all items when
    descriptions=True. Section agendas additionally get a cream rounded
    band with gold border behind the current topic. part_idx/sub_idx
    (legacy call-site signature) map onto the flat M7_OUTLINE index:
    part 0 subs 0-3 -> items 0-3, part 1 subs 0-1 -> items 4-5."""
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, section_tag)
    _draw_action_title(slide, title)

    hi = set()
    if part_idx is not None and sub_idx is not None:
        hi.add(sub_idx + (4 if part_idx == 1 else 0))
    if descriptions:
        hi = set(range(len(M7_OUTLINE)))

    title_h = Inches(0.42)
    desc_h = Inches(0.38)
    gap = Inches(0.11)
    pitch = title_h + desc_h + gap
    total = pitch * len(M7_OUTLINE) - gap
    top = Inches(1.60)
    bottom = Inches(7.02)
    y = int(top + max(0, (bottom - top - total) // 2))

    ys = []
    for i, (item, desc) in enumerate(M7_OUTLINE):
        ys.append(y)
        if not descriptions and i in hi:
            band = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, int(Inches(0.90)),
                int(y - Inches(0.06)), int(Inches(12.15)),
                int(title_h + desc_h + Inches(0.10)))
            band.adjustments[0] = 0.35
            band.fill.solid()
            band.fill.fore_color.rgb = CREAM
            band.line.color.rgb = GOLD
            band.line.width = Pt(1.0)
            band.shadow.inherit = False
            _add_drop_shadow(band)
        circ = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, int(Inches(1.15)), int(y + Inches(0.02)),
            int(Inches(0.58)), int(Inches(0.58)))
        circ.fill.solid()
        circ.fill.fore_color.rgb = GOLD
        circ.line.fill.background()
        circ.shadow.inherit = False
        tf = circ.text_frame
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = str(i + 1)
        run.font.size = Pt(25)
        run.font.bold = True
        run.font.color.rgb = NAVY
        run.font.name = "Calibri"
        rows = [([(item[0].upper() + item[1:],
                   {'bold': True, 'size': 25, 'color': NAVY})], 0,
                 {'bullet_style': 'none', 'space_before_pts': 0})]
        if i in hi:
            rows.append(([(desc, {'size': 22, 'color': GRAY})], 0,
                         {'bullet_style': 'none', 'space_before_pts': 0}))
        _add_hierarchical_bullets(
            slide, Inches(2.05), y, Inches(11.0), title_h + desc_h,
            rows, size=25, line_spacing_pts=0)
        y = int(y + pitch)

    slide._outline_item_ys = ys
    _draw_footer(slide, FOOTER_TEXT, page_num)
    return slide


def make_stub(prs, page_num, section_tag, title, note, *, hidden=False):
    """Positional placeholder so the 87-slide numbering is stable from the
    first build. Replaced by a real builder (or a phase-3 splice) later."""
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, section_tag)
    _draw_action_title(slide, title)
    _add_text(slide, MARGIN, Inches(3.4), RULE_W, Inches(0.8),
              f"[ {note} ]", size=20, italic=True, color=GRAY,
              font="Calibri", align=PP_ALIGN.CENTER)
    _draw_footer(slide, FOOTER_TEXT, page_num)
    if hidden:
        slide._element.set('show', '0')
    return slide


# --------------------------------------------------------------------------
# Front matter
# --------------------------------------------------------------------------

def slide_02_recap(prs):
    bullets = [
        ("Price discrimination: different customers", 0),
        ("First degree (perfect discrimination)", 1),
        ("Third degree: segmentation", 1),
        ("Higher price to less elastic groups", 2, {'size': 20}),
        ("Second degree (indirect): versioning", 1),
        ("Volume pricing: same consumer, multiple units", 0),
        ("Unlimited access fee if MC=0 (e.g., Netflix)", 1),
        ("Fee = all the area under the demand curve", 2, {'size': 20}),
        ("Two-part pricing when MC>0 (e.g., Starbucks)", 1),
        ("Per-unit fee = MC", 2, {'size': 20}),
        ("Flat fee = area under the demand curve above the per-unit fee",
         2, {'size': 20}),
        ("Both can be computed directly from the demand function", 0),
    ]
    s = make_content_bulleted(
        prs, page_num=2, section_tag=TAG_RECAP,
        title="Lessons from Module 6: Complex Pricing",
        bullets=bullets, size=24, sub_size=22, line_spacing_pts=10)
    return s


def slide_03_roadmap(prs):
    """Course roadmap (mirrors Module 3's agenda flowchart) — Module 7
    lives in course part 4: Markets, Pricing and Strategy (navy box)."""

    def draw(slide):
        box_h = Inches(0.85)
        narrow_w = Inches(4.6)
        wide_w = Inches(8.6)
        gap = Inches(0.3)
        slide_mid = SLIDE_W // 2

        top_x = slide_mid - wide_w // 2
        top_y = Inches(1.95)
        _add_rounded_filled_box(slide, top_x, top_y, wide_w, box_h,
                                "1. The Economic Way of Thinking",
                                fill=FADED, text_color=WHITE, size=24,
                                bold=True)

        row2_y = Inches(3.55)
        left_x = slide_mid - gap // 2 - narrow_w
        right_x = slide_mid + gap // 2
        _add_rounded_filled_box(slide, left_x, row2_y, narrow_w, box_h,
                                "2. Value and Demand",
                                fill=FADED, text_color=WHITE, size=26,
                                bold=True)
        _add_rounded_filled_box(slide, right_x, row2_y, narrow_w, box_h,
                                "3. Supply and Cost",
                                fill=FADED, text_color=WHITE, size=26,
                                bold=True)

        bot_x = slide_mid - wide_w // 2
        bot_y = Inches(5.35)
        _add_rounded_filled_box(slide, bot_x, bot_y, wide_w, box_h,
                                "4. Markets, Pricing and Strategy",
                                fill=NAVY, text_color=WHITE, size=24,
                                bold=True)

        top_bottom_y = top_y + box_h
        _add_arrow(slide, (top_x + wide_w // 2, top_bottom_y),
                   (left_x + narrow_w // 2, row2_y),
                   color=FADED, weight_pt=3.0, head=True)
        _add_arrow(slide, (top_x + wide_w // 2, top_bottom_y),
                   (right_x + narrow_w // 2, row2_y),
                   color=FADED, weight_pt=3.0, head=True)
        row2_bottom_y = row2_y + box_h
        _add_arrow(slide, (left_x + narrow_w // 2, row2_bottom_y),
                   (bot_x + wide_w // 2, bot_y),
                   color=FADED, weight_pt=3.0, head=True)
        # gray like the other connectors — Nico 2026-08-03 (was navy 3.5)
        _add_arrow(slide, (right_x + narrow_w // 2, row2_bottom_y),
                   (bot_x + wide_w // 2, bot_y),
                   color=FADED, weight_pt=3.0, head=True)

        # "we are here" — gold label right of box 4, arrow into its edge
        lbl_x = bot_x + wide_w + Inches(0.45)
        lbl_y = bot_y + box_h // 2 - Inches(0.16)
        _add_text(slide, lbl_x, lbl_y, Inches(1.9), Inches(0.32),
                  "we are here", size=16, italic=True, bold=True,
                  color=GOLD, font="Calibri", align=PP_ALIGN.LEFT)
        _add_arrow(slide, (lbl_x - Inches(0.08), lbl_y + Inches(0.16)),
                   (bot_x + wide_w + Inches(0.04), lbl_y + Inches(0.16)),
                   color=GOLD, weight_pt=2.5, head=True)

    s = make_diagram_slide(prs, page_num=3, section_tag=TAG_ROADMAP,
                           title="Agenda for the Class",
                           draw_diagram=draw)
    _set_notes(s, (
        "This is where Module 7 sits in the 405 course as a whole. We built "
        "the economic way of thinking, then demand, then supply and cost. "
        "Module 7 continues part 4: markets, pricing, and strategy — moving "
        "from monopoly and monopolistic competition to oligopoly, where a "
        "few firms interact strategically, and to game theory, the toolkit "
        "for analyzing that interaction."))
    return s


# --------------------------------------------------------------------------
# Part 1 · Oligopoly (5–9)
# --------------------------------------------------------------------------

def slide_05_concept_map(prs):
    """Concept-map overview: how Module 7's pieces connect (added
    2026-08-10, modeled on Module 3's "How the Pieces Connect")."""

    def draw(slide):
        _add_text(slide, Inches(0.3), Inches(1.32), Inches(6.0),
                  Inches(0.42), "OLIGOPOLY  (Part 1)", size=22,
                  bold=True, color=GRAY, font="Calibri",
                  align=PP_ALIGN.CENTER)
        _add_text(slide, Inches(7.03), Inches(1.32), Inches(6.0),
                  Inches(0.42), "GAME THEORY  (Part 2)", size=22,
                  bold=True, color=GRAY, font="Calibri",
                  align=PP_ALIGN.CENTER)

        # roots
        l_x, l_y, l_w, l_h = Inches(1.0), Inches(1.82), Inches(4.6), \
            Inches(0.9)
        _add_rounded_filled_box(slide, l_x, l_y, l_w, l_h,
                        "Oligopoly\nfew large firms — each firm's best "
                        "move depends on its rivals'",
                        fill=NAVY, text_color=WHITE, size=14, bold=True,
                        corner_pct=0.08, shadow=True)
        r_x, r_y, r_w, r_h = Inches(7.73), Inches(1.82), Inches(4.6), \
            Inches(0.9)
        _add_rounded_filled_box(slide, r_x, r_y, r_w, r_h,
                        "Game Theory\nthe toolkit for strategic "
                        "interaction",
                        fill=NAVY, text_color=WHITE, size=14, bold=True,
                        corner_pct=0.08, shadow=True)

        # left children: the three competition modes
        ch_y, ch_h = Inches(3.0), Inches(1.9)
        gap = Inches(0.2)
        cw = (Inches(6.0) - 2 * gap) // 3
        lx1, lx2, lx3 = (Inches(0.3), Inches(0.3) + cw + gap,
                         Inches(0.3) + 2 * (cw + gap))
        _add_rounded_filled_box(slide, lx1, ch_y, cw, ch_h,
                        "Collusion / Cartel\n\nact jointly as a "
                        "monopolist\n\nprofitable, but unstable "
                        "(and illegal)",
                        fill=NAVY, text_color=WHITE, size=13, bold=True,
                        corner_pct=0.08, shadow=True)
        _add_rounded_filled_box(slide, lx2, ch_y, cw, ch_h,
                        "Cournot\n\ncompete on quantities\n\nprice "
                        "between monopoly and competition",
                        fill=NAVY, text_color=WHITE, size=13, bold=True,
                        corner_pct=0.08, shadow=True)
        _add_rounded_filled_box(slide, lx3, ch_y, cw, ch_h,
                        "Bertrand\n\ncompete on prices\n\nidentical goods:\nP = MC\ndifferentiated goods:\nP > MC",
                        fill=NAVY, text_color=WHITE, size=13, bold=True,
                        corner_pct=0.08, shadow=True)

        # right children: the three key concepts
        rx1, rx2, rx3 = (Inches(7.03), Inches(7.03) + cw + gap,
                         Inches(7.03) + 2 * (cw + gap))
        _add_rounded_filled_box(slide, rx1, ch_y, cw, ch_h,
                        "Dominant strategy\n\nbest no matter what the "
                        "rival does\n\n(Prisoner's Dilemma)",
                        fill=NAVY, text_color=WHITE, size=13, bold=True,
                        corner_pct=0.08, shadow=True)
        _add_rounded_filled_box(slide, rx2, ch_y, cw, ch_h,
                        "Nash equilibrium\n\nindividually best responses\n\n"
                        "⇒ no player wants to deviate",
                        fill=NAVY, text_color=WHITE, size=13, bold=True,
                        corner_pct=0.08, shadow=True)
        _add_rounded_filled_box(slide, rx3, ch_y, cw, ch_h,
                        "Commitment\n\nbind yourself to change the "
                        "game\n\n(M.A.D., entry deterrence)",
                        fill=NAVY, text_color=WHITE, size=13, bold=True,
                        corner_pct=0.08, shadow=True)

        # fan-out arrows (roots -> children)
        for cx in (lx1, lx2, lx3):
            _add_arrow(slide, (l_x + l_w // 2, l_y + l_h),
                       (cx + cw // 2, ch_y), color=NAVY,
                       weight_pt=2.5, head=True)
        for cx in (rx1, rx2, rx3):
            _add_arrow(slide, (r_x + r_w // 2, r_y + r_h),
                       (cx + cw // 2, ch_y), color=NAVY,
                       weight_pt=2.5, head=True)

        # MB = MC anchor star under Cournot (each firm optimizes
        # given the rival's action — same star as slides 20/24/25)
        sun_w, sun_h = Inches(1.8), Inches(1.3)
        sun_x = lx2 + (cw - sun_w) // 2
        sun_y = Inches(5.55)
        _add_anchor_burst(slide, sun_x, sun_y, sun_w, sun_h,
                          top_text="MB = MC",
                          bottom_text="(given the rival's action)",
                          top_size=16, bottom_size=11)
        _add_arrow(slide, (sun_x + sun_w // 2, sun_y),
                   (lx2 + cw // 2, ch_y + ch_h), color=GOLD,
                   weight_pt=2.5, head=True)

        # bridge: oligopoly models ARE games -> Nash equilibrium
        b_x, b_y = Inches(6.55), Inches(5.65)
        b_w, b_h = Inches(3.6), Inches(0.95)
        _add_outlined_box(slide, b_x, b_y, b_w, b_h,
                          "Cournot & Bertrand are games:\ntheir "
                          "equilibria are Nash equilibria",
                          fill=WHITE, line=GOLD, text_color=NAVY,
                          size=14, bold=True, line_w=2.0,
                          rounded=True, shadow=True)
        _add_arrow(slide, (lx3 + cw - Inches(0.3), ch_y + ch_h),
                   (b_x + Inches(0.5), b_y), color=GOLD,
                   weight_pt=3.0, head=True)
        # second inflow from the Cournot side (Nico's hand edit,
        # 2026-08-11)
        _add_arrow(slide, (Inches(3.753), Inches(4.952)),
                   (Inches(6.55), Inches(6.125)), color=GOLD,
                   weight_pt=3.0, head=True)
        _add_arrow(slide, (b_x + b_w - Inches(0.6), b_y),
                   (rx2 + cw // 2, ch_y + ch_h), color=GOLD,
                   weight_pt=3.0, head=True)

    s = make_diagram_slide(
        prs, page_num=5,
        section_tag="Module 7 · Concept Map",
        title="How the Pieces of Module 7 Connect",
        draw_diagram=draw,
    )
    _set_notes(s, (
        "The map of Module 7 before we dive in. Part 1 is oligopoly — "
        "few large firms whose best moves depend on each other — in "
        "three flavors: collusion (act jointly as a monopolist, "
        "profitable but unstable and illegal), Cournot competition on "
        "quantities, and Bertrand competition on prices. Part 2 is "
        "game theory, the general toolkit behind all of this: dominant "
        "strategies, Nash equilibrium, and commitment. The gold bridge "
        "is the punchline connecting the two halves: Cournot and "
        "Bertrand are themselves games, and their equilibria are Nash "
        "equilibria. We return to this map as the module unfolds."))
    return s


def slide_06_market_structures(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_OLIGO)
    _draw_action_title(slide, "The Four Basic Market Structures")

    rows = [
        ["", "Perfect competition", "Monopolistic competition",
         "Oligopoly", "Monopoly"],
        ["Number of Firms", "Many", "Many", "Few", "One"],
        ["Type of Products Sold", "Identical", "Differentiated",
         "Identical or Differentiated", "Unique"],
        ["Barriers to Entry", "None", "None", "Some", "Many"],
    ]
    tbl_w, tbl_h = Inches(12.2), Inches(3.9)
    left = (SLIDE_W - tbl_w) // 2
    top = Inches(2.15)
    col_w = [Inches(2.6)] + [Inches(2.4)] * 4
    _add_styled_table(
        slide, left, top, tbl_w, tbl_h, rows,
        col_widths=col_w,
        row_heights=[Inches(0.9), Inches(0.9), Inches(1.2), Inches(0.9)],
        font_size=20, header_size=20)
    # red circle around the Oligopoly header — like the original slide
    # (replaces the earlier gold header cell — Nico 2026-08-03)
    oval = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left + Inches(2.6) + Inches(2.4) * 2 - Inches(0.12),
        top - Inches(0.13), Inches(2.64), Inches(1.16))
    oval.fill.background()
    oval.line.color.rgb = RED
    oval.line.width = Pt(2.5)
    oval.shadow.inherit = False
    # (bottom italic note deleted by hand — Nico 2026-08-03)
    _draw_footer(slide, FOOTER_TEXT, 6)
    return slide


RED = RGBColor(0xC0, 0x00, 0x00)


def _double_arrow(slide, start_xy, end_xy, *, color=NAVY, weight_pt=3.0,
                  head_size='med'):
    """Straight connector with triangle arrowheads on BOTH ends."""
    shp = _add_arrow(slide, start_xy, end_xy, color=color,
                     weight_pt=weight_pt, head=True, head_size=head_size)
    ln = shp.line._get_or_add_ln()
    tail = ln.find(qn('a:tailEnd'))
    head = ln.makeelement(qn('a:headEnd'), {'type': 'triangle',
                                            'w': head_size,
                                            'h': head_size})
    tail.addprevious(head)
    return shp


ACC_ORANGE = RGBColor(0xED, 0x7D, 0x31)   # accent6 orange (Nico, slide 6)


def _capped_line(slide, start_xy, end_xy, *, color=ACC_ORANGE,
                 weight_pt=2.75):
    """Straight connector with small OVAL caps on both ends."""
    shp = _add_arrow(slide, start_xy, end_xy, color=color,
                     weight_pt=weight_pt, head=False)
    ln = shp.line._get_or_add_ln()
    for tag in ('a:headEnd', 'a:tailEnd'):
        ln.append(ln.makeelement(qn(tag), {'type': 'oval'}))
    return shp


def slide_07_spectrum(prs):
    """Market-power spectrum. Geometry = Nico's hand-edited layout, ported
    verbatim 2026-08-03 (second pass): 28 pt end labels; orange ribbon
    bars + oval-capped zone lines for (Price Takers)/(Price Searchers);
    slim red arrow + 'WE ARE HERE' label + tightened red Oligopoly box
    (grouped by hand — grouping re-done in phase 3); adidas replaced by
    the 'Restaurants in WeHo' map with caption."""
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_OLIGO)
    _draw_action_title(slide, "Where We Are and Where We Are Going")

    # "Least / Most Market Power" labels above the spectrum arrow
    _add_text(slide, Inches(0.550), Inches(1.780), Inches(3.60),
              Inches(0.471), "Least Market Power", size=28, bold=True,
              color=NAVY, font="Calibri", align=PP_ALIGN.LEFT)
    _add_text(slide, Inches(9.210), Inches(1.780), Inches(3.59),
              Inches(0.471), "Most Market Power", size=28, bold=True,
              color=NAVY, font="Calibri", align=PP_ALIGN.RIGHT)

    # main spectrum arrow — double-headed
    _double_arrow(slide, (Inches(0.5), Inches(2.46)),
                  (Inches(12.83), Inches(2.46)),
                  color=NAVY, weight_pt=3.5, head_size='lg')

    # the four structures, separated by red divider bars
    labels = [(2.05, ("Perfect", "Competition")),
              (5.0, ("Monopolistic", "Competition")),
              (8.35, ("Oligopoly",)), (11.3, ("Monopoly",))]
    for cx, lines in labels:
        base = 2.68 if len(lines) == 2 else 2.88
        for li, line in enumerate(lines):
            _add_text(slide, Inches(cx - 1.55), Inches(base + 0.40 * li),
                      Inches(3.1), Inches(0.42), line,
                      size=22, bold=True, color=NAVY, font="Calibri",
                      align=PP_ALIGN.CENTER)
    for dx in [3.53, 6.67, 9.83]:
        _add_rect(slide, Inches(dx), Inches(2.66), Inches(0.045),
                  Inches(0.85), RED)

    # (Price Takers)/(Price Searchers): rounded, shadowed orange boxes
    # (zone lines removed — Nico 2026-08-04)
    for lx, w, txt in [(1.071, 2.121, "(Price Takers)"),
                       (4.330, 7.946, "(Price Searchers)")]:
        shp = _add_rounded_filled_box(
            slide, Inches(lx), Inches(3.879), Inches(w), Inches(0.404),
            txt, fill=ACC_ORANGE, text_color=WHITE, size=24, bold=False,
            corner_pct=0.35, shadow=True)
        shp.text_frame.paragraphs[0].runs[0].font.italic = True

    # pictures (flat — logo exception), hand-tuned positions
    _add_media_image(slide, "image14.jpeg", left=Inches(1.273),
                     top=Inches(4.720), width=Inches(1.554),
                     height=Inches(1.500), rounded=False, shadow=False)
    _add_media_image(slide, "_s6_image6.png", left=Inches(3.988),
                     top=Inches(4.502), width=Inches(2.387),
                     height=Inches(2.267), rounded=False, shadow=False)
    _add_text(slide, Inches(4.107), Inches(6.721), Inches(2.091),
              Inches(0.376), "Restaurants in WeHo", size=16,
              color=RGBColor(0, 0, 0), font="Calibri",
              align=PP_ALIGN.CENTER)
    _add_media_image(slide, "image15.jpeg", left=Inches(7.786),
                     top=Inches(4.720), width=Inches(1.128),
                     height=Inches(0.750), rounded=False, shadow=False)
    _add_media_image(slide, "image16.png", left=Inches(7.875),
                     top=Inches(5.620), width=Inches(0.949),
                     height=Inches(0.750), rounded=False, shadow=False)
    _add_media_image(slide, "image17.jpeg", left=Inches(10.591),
                     top=Inches(4.720), width=Inches(1.417),
                     height=Inches(1.500), rounded=False, shadow=False)

    # "WE ARE HERE" cluster: slim red arrow into the red Oligopoly box
    # — grouped into ONE object (Nico's hand group, 2026-08-12)
    wah = _add_text(slide, Inches(5.505), Inches(1.493), Inches(1.620),
                    Inches(0.370), "WE ARE HERE", size=16, bold=True,
                    color=RED, font="Calibri", align=PP_ALIGN.LEFT)
    ra = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(5.200),
                                Inches(1.770), Inches(2.350),
                                Inches(0.250))
    ra.fill.solid()
    ra.fill.fore_color.rgb = RED
    ra.line.fill.background()
    ra.shadow.inherit = False
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.535),
                                 Inches(1.620), Inches(1.675),
                                 Inches(5.050))
    box.fill.background()
    box.line.color.rgb = RED
    box.line.width = Pt(3.5)
    box.shadow.inherit = False
    _group_shapes(slide, [wah, ra, box], "WeAreHere")

    _draw_footer(slide, FOOTER_TEXT, 7)
    return slide


def slide_08_characteristics(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_OLIGO)
    _draw_action_title(slide, "Market Structure Characteristics")

    rows = [
        ["", "Oligopoly"],
        ["Number of actual or potential competitors",
         "Few enough sellers so that their decisions are "],
        ["Product differentiation", "High or low"],
        ["Entry conditions",
         "High barriers due to economies of scale, capital requirements, "
         "access"],
        ["Profit potential",
         "Economic profit can be positive in short/long run"],
        ["Examples",
         "Cell phone providers; soft drinks; tobacco; semiconductors; "
         "wide-bodied aircraft"],
    ]
    tbl_w, tbl_h = Inches(11.8), Inches(4.7)
    left = (SLIDE_W - tbl_w) // 2
    top = Inches(1.8)
    gf = _add_styled_table(
        slide, left, top, tbl_w, tbl_h, rows,
        col_widths=[Inches(4.1), Inches(7.7)],
        row_heights=[Inches(0.55)] + [Inches(0.83)] * 5,
        font_size=18, header_size=20)
    # "interdependent" in red — hand-edit ported (Nico 2026-08-04)
    p = gf.table.cell(1, 1).text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "interdependent"
    run.font.name = "Calibri"
    run.font.size = Pt(18)
    run.font.color.rgb = RED
    _draw_footer(slide, FOOTER_TEXT, 8)
    _set_notes(slide, (
        "Oligopoly means markets with few sellers, but more than one "
        "(which is monopoly). How few? What is the distinction between "
        "oligopoly and monopolistic competition?\n"
        "There are barriers to entry in oligopolistic markets, keeping the "
        "number of competitors low and creating the potential for long "
        "term abnormal profits.\n"
        "There is strategic interdependence in oligopoly: under "
        "monopolistic competition, each firm ignores price changes by "
        "other firms because each is so small their price behavior does "
        "not affect your behavior. Under oligopoly, each has an incentive "
        "to respond to others' deviations.\n"
        "Examples of oligopoly include tobacco, satellite radio (before "
        "the merger), wide-bodied aircraft, and mobile phone services."))
    return slide


def slide_09_concentration(prs):
    """Concentration example UPDATED 2026-07-29 (research-verified):
    top-3 US wireless ≈ 97% of retail subscribers; T-Mobile absorbed
    Sprint (2020) and UScellular's wireless operations (Aug 2025).
    Subscriber counts: Verizon ~146M, T-Mobile ~140M, AT&T ~119M
    (company reports, late 2025)."""
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_OLIGO)
    _draw_action_title(slide, "Measuring Market Concentration")

    bullets = [
        ("When do we have an oligopoly?", 0),
        ("Simple measure: Concentration Ratio CR4", 0),
        ("Market share of the four largest firms", 1),
        ("Oligopoly if CR4 > 66%", 1),
        ("Example: US wireless carriers — top 3 ≈ 97%", 1),
        ("Limitations of these measures", 0),
        ("Definition of relevant market", 1),
        # italic "potential" — hand-edit ported (Nico 2026-08-04)
        ([("Role of ", {}), ("potential", {'italic': True}),
          (" competition", {})], 1),
    ]
    # bullets bumped 24/22 → 28/24 by hand (Nico 2026-08-04)
    _add_hierarchical_bullets(slide, left=MARGIN, top=Inches(1.75),
                              width=Inches(7.1), height=Inches(4.0),
                              items=bullets, size=28, sub_size=24,
                              line_spacing_pts=12)

    # native bar mini-chart, US wireless subscribers late 2025
    # (moved up 0.6" by Nico, 2026-08-12; whole panel = one group)
    _s9_before = _shape_ids(slide)
    chart_l, chart_t = Inches(8.15), Inches(1.65)
    chart_w, chart_h = Inches(4.75), Inches(4.05)
    _add_graphicframe_shadow(slide, chart_l, chart_t, chart_w, chart_h)
    _add_text(slide, chart_l, chart_t + Inches(0.12), chart_w, Inches(0.55),
              "US wireless subscribers,\nlate 2025 (millions)",
              size=15, italic=True, bold=True, color=NAVY, font="Calibri",
              align=PP_ALIGN.CENTER)
    # brand-colored bars, each with a small logo chip (Nico 2026-08-03).
    # Logos: Wikimedia Commons (_logo_*.png in _source_images).
    ASPECTS = {"_logo_verizon.png": 500 / 111, "_logo_tmobile.png": 1.0,
               "_logo_att.png": 500 / 206}
    vals = [("Verizon", 146, RGBColor(0xEE, 0x00, 0x00),
             "_logo_verizon.png"),
            ("T-Mobile", 140, RGBColor(0xE2, 0x00, 0x74),
             "_logo_tmobile.png"),
            ("AT&T", 119, RGBColor(0x00, 0xA8, 0xE0), "_logo_att.png")]
    base_y = chart_t + Inches(3.05)
    max_h = Inches(1.85)
    bar_w = Inches(1.05)
    gap = Inches(0.38)
    x = chart_l + Inches(0.45)
    for name, v, bcolor, logo in vals:
        h = int(max_h * v / 146)
        _add_rect(slide, x, base_y - h, bar_w, h, bcolor)
        _add_text(slide, x - Inches(0.2), base_y - h - Inches(0.36),
                  bar_w + Inches(0.4), Inches(0.32), str(v),
                  size=16, bold=True, color=NAVY, font="Calibri",
                  align=PP_ALIGN.CENTER)
        _add_text(slide, x - Inches(0.2), base_y + Inches(0.08),
                  bar_w + Inches(0.4), Inches(0.32), name,
                  size=15, color=NAVY, font="Calibri",
                  align=PP_ALIGN.CENTER)
        # white chip inside the bar carrying the company logo
        chip_w, chip_h = Inches(0.92), Inches(0.42)
        chip = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            int(x + (bar_w - chip_w) // 2), int(base_y - Inches(0.55)),
            int(chip_w), int(chip_h))
        chip.adjustments[0] = 0.30
        chip.fill.solid()
        chip.fill.fore_color.rgb = WHITE
        chip.line.fill.background()
        chip.shadow.inherit = False
        asp = ASPECTS[logo]
        lh = min(0.30, 0.80 / asp)
        lw = lh * asp
        _add_media_image(
            slide, logo,
            left=int(x + bar_w // 2 - Inches(lw) // 2),
            top=int(base_y - Inches(0.55) + (chip_h - Inches(lh)) // 2),
            width=Inches(lw), height=Inches(lh),
            rounded=False, shadow=False)
        x += bar_w + gap
    _add_rect(slide, chart_l + Inches(0.3), base_y, chart_w - Inches(0.6),
              Inches(0.018), NAVY)
    _add_text(slide, chart_l, base_y + Inches(0.42), chart_w, Inches(0.3),
              "Source: company reports (retail connections)",
              size=11, italic=True, color=GRAY, font="Calibri",
              align=PP_ALIGN.CENTER)
    _group_shapes(slide, _new_shapes_since(slide, _s9_before),
                  "WirelessChart")

    # merger annotation — cream card, rebuilt at Nico's hand position
    # bottom-right under the chart (2026-08-12); card + text + logos
    # = ONE group
    _s9_before2 = _shape_ids(slide)
    mb_l, mb_t = Inches(6.666), Inches(6.093)
    mb_w, mb_h = Inches(6.391), Inches(0.838)
    mbox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  int(mb_l), int(mb_t), int(mb_w),
                                  int(mb_h))
    mbox.adjustments[0] = 0.12
    mbox.fill.solid()
    mbox.fill.fore_color.rgb = CREAM
    mbox.line.color.rgb = NAVY
    mbox.line.width = Pt(1.0)
    mbox.shadow.inherit = False
    tb = slide.shapes.add_textbox(int(Inches(6.886)), int(Inches(6.12)),
                                  int(Inches(5.995)), int(Inches(0.35)))
    p = tb.text_frame.paragraphs[0]
    for txt, bold in [("Mergers:  ", True),
                      ("T-Mobile absorbed Sprint (2020) and UScellular "
                       "(2025)", False)]:
        r = p.add_run()
        r.text = txt
        r.font.name = "Calibri"
        r.font.size = Pt(16)
        r.font.bold = bold
        r.font.color.rgb = NAVY
    for fname, lx, ty, w_in, h_in in [
            ("_logo_tmobile.png", 8.02, 6.565, 0.27, 0.30),
            ("_logo_sprint.png", 8.578, 6.575, 0.607, 0.28),
            ("_logo_tmobile.png", 9.725, 6.565, 0.27, 0.30),
            ("_logo_uscellular.png", 10.283, 6.605, 0.967, 0.22)]:
        _add_media_image(slide, fname, left=Inches(lx), top=Inches(ty),
                         width=Inches(w_in), height=Inches(h_in),
                         rounded=False, shadow=False)
    for px in (8.308, 10.013):
        _add_text(slide, Inches(px), Inches(6.546), Inches(0.234),
                  Inches(0.34), "+", size=18, bold=True, color=NAVY,
                  font="Calibri", align=PP_ALIGN.CENTER)
    _group_shapes(slide, _new_shapes_since(slide, _s9_before2),
                  "MergerCard")

    _draw_footer(slide, FOOTER_TEXT, 9)
    _set_notes(slide, (
        "Example: concentration in US mobile phone services — clearly an "
        "oligopoly. Verizon, T-Mobile, and AT&T together serve roughly 97% "
        "of retail wireless subscribers. And consolidation continues: "
        "T-Mobile absorbed Sprint in 2020 and UScellular's wireless "
        "operations in August 2025."))
    return slide


# --------------------------------------------------------------------------
# §1.1 Collusion and Cartels (11)
# --------------------------------------------------------------------------

def slide_12_collusion(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_COLLU)
    _draw_action_title(slide, "Collusion")

    box_w = Inches(11.4)
    left = (SLIDE_W - box_w) // 2

    # Outer frame titled "Definition of Collusion" (label at Nico's
    # hand-set 24 pt gold, 2026-08-04) with the navy hero box NESTED
    # inside. Frame + label + hero = ONE group (Nico 2026-08-12).
    _s12_a = _shape_ids(slide)
    _add_outlined_box(
        slide, left, Inches(1.42), box_w, Inches(1.90), "",
        line=NAVY, fill=WHITE, line_w=1.25, rounded=True, shadow=True,
        corner_pct=0.08)
    _add_text(slide, left + Inches(0.30), Inches(1.54), Inches(6.0),
              Inches(0.45), "Definition of Collusion", size=24, bold=True,
              color=GOLD, font="Calibri")
    _add_rounded_filled_box(
        slide, left + Inches(0.25), Inches(2.10), box_w - Inches(0.50),
        Inches(1.00),
        "Firms in an oligopoly coordinate their production and pricing "
        "decisions",
        fill=NAVY, text_color=WHITE, size=24, bold=True, corner_pct=0.08)
    _group_shapes(slide, _new_shapes_since(slide, _s12_a), "DefCard")

    _s12_b = _shape_ids(slide)
    _add_convention_box(
        slide, left, Inches(3.55), box_w, Inches(1.8),
        runs=[
            ("▪  Aim: collectively act as a monopolist to gain monopoly "
             "profits", {'size': 22, 'color': NAVY}),
            ("▪  Profits to be split among the participating firms",
             {'size': 22, 'color': NAVY, 'newline': True}),
            ("▪  Extreme case – a cartel: an explicit, organized "
             "price-fixing agreement (e.g., OPEC)",
             {'size': 22, 'color': NAVY, 'newline': True}),
        ],
        size=22)
    _group_shapes(slide, _new_shapes_since(slide, _s12_b), "AimCard")

    _s12_c = _shape_ids(slide)
    _add_convention_box(
        slide, left, Inches(5.60), box_w, Inches(1.40),
        runs=[
            ("“People of the same trade seldom meet together, even for "
             "merriment and diversion, but … in some contrivance to raise "
             "prices.”", {'size': 20, 'italic': True, 'color': NAVY}),
            ("— Adam Smith (1723–1790)",
             {'size': 17, 'bold': True, 'color': GRAY, 'newline': True}),
        ],
        fill_rgb=WHITE, border=GOLD, line_w=1.5, size=20,
        align=PP_ALIGN.CENTER)
    _group_shapes(slide, _new_shapes_since(slide, _s12_c), "QuoteCard")

    _draw_footer(slide, FOOTER_TEXT, 12)
    return slide


# --------------------------------------------------------------------------
# §1.2 Cournot (17–21)
# --------------------------------------------------------------------------

def slide_18_note(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_COURNOT)
    _draw_action_title(slide, "Note")

    box_w = Inches(9.0)
    left = (SLIDE_W - box_w) // 2
    _add_convention_box(
        slide, left, Inches(2.8), box_w, Inches(2.4),
        runs=[
            ("From now on:", {'size': 26, 'bold': True, 'color': NAVY}),
            ("No collusion among competitors",
             {'size': 40, 'color': NAVY, 'newline': True}),
            ("→  Firms compete strategically (“games”)",
             {'size': 24, 'bold': True, 'color': NAVY, 'newline': True}),
        ],
        size=24)
    _draw_footer(slide, FOOTER_TEXT, 18)
    return slide


def slide_19_two_models(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_COURNOT)
    _draw_action_title(slide, "Two Models of Oligopoly Competition")

    col_w = Inches(4.8)
    centers = [Inches(3.55), Inches(9.78)]
    # portrait geometry hand-tuned by Nico (2026-08-04)
    names = [
        ("Cournot (1801–1877)", "image27.png",
         "Cournot Model: Competition in Quantity",
         (2.48, 2.23, 2.23, 3.14)),
        ("Bertrand (1822–1900)", "image26.jpg",
         "Bertrand Model: Competition in Price",
         (8.63, 2.25, 2.59, 3.14)),
    ]
    for cx, (header, img, model, (px, py, pw, ph)) in zip(centers, names):
        _add_text(slide, cx - col_w // 2, Inches(1.75), col_w, Inches(0.45),
                  header, size=24, bold=True, color=NAVY, font="Calibri",
                  align=PP_ALIGN.CENTER)
        _add_media_image(slide, img, left=Inches(px), top=Inches(py),
                         width=Inches(pw), height=Inches(ph),
                         rounded=True, shadow=True)
        _add_rounded_filled_box(
            slide, cx - col_w // 2, Inches(5.6), col_w, Inches(0.85),
            model, fill=NAVY, text_color=WHITE, size=20, bold=True)

    _draw_footer(slide, FOOTER_TEXT, 19)
    _set_notes(slide, (
        "Oligopoly: Cournot versus Bertrand. In what follows we will "
        "assume no collusion: competition rather than cooperation between "
        "firms. There are a lot of models of oligopoly – as many as "
        "assumptions you are willing to make on rivals' responses. You "
        "will come back to these issues in your strategy class, for now "
        "we give the foundations that will be used in strategy.\n"
        "Two oligopoly models – the two most famous ones.\n"
        "Cournot assumed quantity competition: each firm sets quantity "
        "taking as fixed the quantity set by the other. The assumption is "
        "that the other firm will continue to produce the same quantity "
        "whatever I do. As we will see shortly price does not go to the "
        "competitive level but settles between the monopoly and perfect "
        "competition price.\n"
        "Bertrand (in his review of the Cournot book) proposed an "
        "alternative – that competition occurs on prices, not quantity. "
        "That is, the rival's price is held fixed when setting price. The "
        "idea is that in response to my price cut, others don't match. "
        "The equilibrium is where this assumption turns out to be correct "
        "(as before the assumption is false for all except the "
        "equilibrium point). What happens? As we will see shortly, "
        "equilibrium ends up equal to perfect competition."))
    return slide


def slide_20_two_models_bullets(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_COURNOT)
    _draw_action_title(slide, "Two Oligopoly Models")

    col_w = Inches(5.95)
    gap = Inches(0.45)
    left1 = (SLIDE_W - 2 * col_w - gap) // 2
    left2 = left1 + col_w + gap
    hdr_y = Inches(1.85)
    _add_rounded_filled_box(slide, left1, hdr_y, col_w, Inches(0.6),
                            "Cournot", fill=NAVY, text_color=WHITE,
                            size=24, bold=True, corner_pct=0.15)
    _add_rounded_filled_box(slide, left2, hdr_y, col_w, Inches(0.6),
                            "Bertrand", fill=NAVY, text_color=WHITE,
                            size=24, bold=True, corner_pct=0.15)

    cournot = [
        "The rival's supply is taken as given",
        "Firm reacts by adjusting its own quantity accordingly",
        "Price is determined in the market",
    ]
    bertrand = [
        "The rival's price is taken as given",
        "Firm reacts by adjusting its own price (underbidding, if "
        "feasible)",
        "Firms then supply the quantity that is demanded at the price "
        "they set",
    ]
    # shaded cards behind the two bullet columns (Nico 2026-08-04);
    # card + bullets = one unit → group in phase 3
    for lx in (left1, left2):
        _add_outlined_box(slide, lx, Inches(2.75), col_w, Inches(3.55),
                          "", line=NAVY, fill=WHITE, line_w=1.0,
                          rounded=True, shadow=True, corner_pct=0.05)
    _add_bulleted_list(slide, left=left1 + Inches(0.3), top=Inches(3.0),
                       width=col_w - Inches(0.6), height=Inches(3.1),
                       items=cournot, size=22, line_spacing_pts=16)
    _add_bulleted_list(slide, left=left2 + Inches(0.3), top=Inches(3.0),
                       width=col_w - Inches(0.6), height=Inches(3.1),
                       items=bertrand, size=22, line_spacing_pts=16)

    _draw_footer(slide, FOOTER_TEXT, 20)
    return slide


def slide_22_cournot_assumptions(prs):
    """Nico's hand design (2026-08-06): full-bleed background photo,
    assumptions card as his scaled group (injected verbatim; text keeps
    his in-group scaling and spacing)."""
    slide = _blank_slide(prs)
    slide.shapes.add_picture(str(SRC_IMG_DIR / "_s21_image27.png"),
                             3658, 0, width=Inches(13.326),
                             height=Inches(7.5))
    _draw_top_bar_tc(slide, TAG_COURNOT)
    _draw_action_title(slide,
                       "Cournot Competition (with Homogeneous Goods)")
    _inject_handoff_group(slide, "_handoff_s21_group.xml", id_base=9500)
    _draw_footer(slide, FOOTER_TEXT, 22)
    return slide


# --------------------------------------------------------------------------

class SimpleFig:
    """Logical→slide coordinate transform for shape-built charts.
    All geometry in float inches; returns Inches() EMU ints."""

    def __init__(self, left_in, bottom_in, w_in, h_in, xmax, ymax):
        self.l, self.b, self.w, self.h = left_in, bottom_in, w_in, h_in
        self.xmax, self.ymax = xmax, ymax

    def x(self, xv):
        return Inches(self.l + self.w * xv / self.xmax)

    def y(self, yv):
        return Inches(self.b - self.h * yv / self.ymax)


def _fig_axes(slide, fig, *, weight_pt=2.0):
    """Navy axes with triangle arrowheads. Titles are added by the caller
    (they often need OMML subscripts)."""
    _add_arrow(slide, (fig.x(0), fig.y(0)),
               (fig.x(0), Inches(fig.b - fig.h - 0.18)),
               color=NAVY, weight_pt=weight_pt, head=True)
    _add_arrow(slide, (fig.x(0), fig.y(0)),
               (Inches(fig.l + fig.w + 0.18), fig.y(0)),
               color=NAVY, weight_pt=weight_pt, head=True)


def _fig_ytick(slide, fig, val, label, *, color=NAVY, size=16, bold=False):
    return _add_text(slide, Inches(fig.l - 1.07),
                     fig.y(val) - Inches(0.14),
                     Inches(0.95), Inches(0.3), label, size=size,
                     bold=bold, color=color, font="Calibri",
                     align=PP_ALIGN.RIGHT)


def _fig_xtick(slide, fig, val, label, *, color=NAVY, size=16, bold=False):
    return _add_text(slide, fig.x(val) - Inches(0.5),
                     Inches(fig.b + 0.06),
                     Inches(1.0), Inches(0.3), label, size=size,
                     bold=bold, color=color, font="Calibri",
                     align=PP_ALIGN.CENTER)


def _fig_point(slide, fig, xv, yv, *, r_in=0.055, fill=NAVY, line=None):
    d = Inches(2 * r_in)
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 fig.x(xv) - Inches(r_in),
                                 fig.y(yv) - Inches(r_in), d, d)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1.25)
    shp.shadow.inherit = False
    return shp


# OMML shortcuts for firm/market subscripted symbols (upright subscripts —
# they are labels, not variables)
def _oQ(s):
    return _omml_sub(_omml_run('Q'), _omml_text(s))


def _oP(s):
    return _omml_sub(_omml_run('P'), _omml_text(s))


def _oMR(s):
    return _omml_sub(_omml_text('MR'), _omml_text(s))


# --------------------------------------------------------------------------
# Payoff matrix (game-theory conventions: col player on top in accent
# color, row player rotated at left in concept blue, white cells with navy
# borders, payoffs in the players' colors, Nash oval)
# --------------------------------------------------------------------------

ROW_BLUE = RGBColor(0x00, 0x70, 0xC0)


def _payoff_cell(slide, x, y, w, h, a, b, *, row_color, col_color,
                 size=24, text_override=None, dim=None):
    """One payoff cell as a SINGLE shape (text lives inside the rect, so
    no box+textbox grouping is needed). dim ∈ {None,'row','col'} grays
    the other player's payoff (used on the strategy-analysis slides)."""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(x), int(y),
                                 int(w), int(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = WHITE
    shp.line.color.rgb = NAVY
    shp.line.width = Pt(1.5)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    if text_override is not None:
        # caption sits in the UPPER half of the cell; the red cross
        # goes in the lower half, below the text (original layout)
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_top = Inches(0.08)
        r = p.add_run()
        r.text = text_override
        r.font.name = "Calibri"
        r.font.size = Pt(12)   # original caption size; 2-line wrap
        #   keeps the lower half of the cell clear for the red X
        r.font.italic = True
        r.font.color.rgb = GRAY
        return shp
    # dim: the other player's payoff is COVERED by a gray box (see
    # _add_payoff_matrix); render it white so nothing peeks out
    runs = [(str(a), row_color if dim != 'col' else WHITE),
            (" , ", NAVY),
            (str(b), col_color if dim != 'row' else WHITE)]
    for txt, clr in runs:
        r = p.add_run()
        r.text = txt
        r.font.name = "Calibri"
        r.font.size = Pt(size)
        r.font.bold = True
        r.font.color.rgb = clr
    return shp


def _add_payoff_matrix(slide, *, left, top, cell_w, cell_h,
                       row_player, col_player, row_strats, col_strats,
                       payoffs, row_color=ROW_BLUE, col_color=GOLD,
                       payoff_size=24, strat_size=18, name_size=22,
                       caption=None, nash_cells=(), cell_texts=None,
                       dim=None):
    """2×2 (or n×m) payoff matrix. payoffs[r][c] = (a, b); a cell may
    instead take a literal string via cell_texts[(r,c)] (e.g. the
    "impossible due to M.A.D." cells). Returns anchor geometry for
    best-response overlays."""
    left, top = int(left), int(top)
    cell_w, cell_h = int(cell_w), int(cell_h)
    n_r, n_c = len(row_strats), len(col_strats)
    grid_w, grid_h = cell_w * n_c, cell_h * n_r

    # column player name + column strategy labels
    _add_text(slide, left, top - Inches(0.98), grid_w, Inches(0.42),
              col_player, size=name_size, bold=True, color=col_color,
              font="Calibri", align=PP_ALIGN.CENTER)
    for c, s in enumerate(col_strats):
        _add_text(slide, left + c * cell_w, top - Inches(0.56),
                  cell_w, Inches(0.5), s, size=strat_size, bold=True,
                  color=NAVY, font="Calibri", align=PP_ALIGN.CENTER,
                  anchor=MSO_ANCHOR.BOTTOM)

    # row player name (rotated 270°) + row strategy labels
    name_w = grid_h
    cx = left - Inches(2.42)
    cy = top + grid_h // 2
    box = _add_text(slide, int(cx - name_w // 2), int(cy - Inches(0.21)),
                    name_w, Inches(0.42), row_player, size=name_size,
                    bold=True, color=row_color, font="Calibri",
                    align=PP_ALIGN.CENTER)
    box.rotation = 270
    for r, s in enumerate(row_strats):
        _add_text(slide, left - Inches(2.02),
                  int(top + r * cell_h + (cell_h - Inches(0.75)) // 2),
                  Inches(1.9), Inches(0.75), s, size=strat_size, bold=True,
                  color=NAVY, font="Calibri", align=PP_ALIGN.RIGHT,
                  anchor=MSO_ANCHOR.MIDDLE)

    # cells
    anchors = {}
    for r in range(n_r):
        for c in range(n_c):
            x = left + c * cell_w
            y = top + r * cell_h
            override = cell_texts.get((r, c)) if cell_texts else None
            if override is not None:
                _payoff_cell(slide, x, y, cell_w, cell_h, None, None,
                             row_color=row_color, col_color=col_color,
                             size=payoff_size, text_override=override)
            else:
                a, b = payoffs[r][c]
                _payoff_cell(slide, x, y, cell_w, cell_h, a, b,
                             row_color=row_color, col_color=col_color,
                             size=payoff_size, dim=dim)
            anchors[(r, c, 'row')] = (x + int(cell_w * 0.34),
                                      y + cell_h // 2)
            anchors[(r, c, 'col')] = (x + int(cell_w * 0.66),
                                      y + cell_h // 2)
            anchors[(r, c, 'cell')] = (x + cell_w // 2, y + cell_h // 2)
            if override is None and dim in ('row', 'col'):
                # ORIGINAL-deck technique: cover the non-relevant
                # payoff with a small gray box (Nico 2026-08-09)
                bx, by = anchors[(r, c,
                                  'col' if dim == 'row' else 'row')]
                cover = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, int(bx - Inches(0.23)),
                    int(by - Inches(0.17)), int(Inches(0.46)),
                    int(Inches(0.34)))
                cover.fill.solid()
                cover.fill.fore_color.rgb = GRAY
                cover.line.fill.background()
                cover.shadow.inherit = False

    # Nash oval(s) — the finale element
    for (r, c) in nash_cells:
        ow, oh = int(cell_w + Inches(0.22)), int(cell_h + Inches(0.20))
        cxx, cyy = anchors[(r, c, 'cell')]
        oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, cxx - ow // 2,
                                      cyy - oh // 2, ow, oh)
        oval.fill.background()
        oval.line.color.rgb = GOLD
        oval.line.width = Pt(2.75)
        oval.shadow.inherit = False

    # caption ("Payoffs to (Row, Col)") below the grid
    if caption:
        tb = slide.shapes.add_textbox(left - Inches(0.4),
                                      top + grid_h + Inches(0.14),
                                      grid_w + Inches(0.8), Inches(0.35))
        tf = tb.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        for txt, clr, bold in [("Payoffs to (", NAVY, False),
                               (row_player, row_color, True),
                               (", ", NAVY, False),
                               (col_player, col_color, True),
                               (")", NAVY, False)]:
            r_ = p.add_run()
            r_.text = txt
            r_.font.name = "Calibri"
            r_.font.size = Pt(15)
            r_.font.italic = True
            r_.font.bold = bold
            r_.font.color.rgb = clr
    return anchors


def _br_circle(slide, anchor, color, *, w_in=0.62, h_in=0.46,
               weight_pt=2.25):
    """Best-response circle around one payoff number."""
    cx, cy = anchor
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 int(cx - Inches(w_in / 2)),
                                 int(cy - Inches(h_in / 2)),
                                 Inches(w_in), Inches(h_in))
    shp.fill.background()
    shp.line.color.rgb = color
    shp.line.width = Pt(weight_pt)
    shp.shadow.inherit = False
    return shp


# --------------------------------------------------------------------------
# Slide 9 — nickel example (UPDATED 2026-07-29: Vale + Tsingshan as the
# stylized duopoly; Norilsk excluded per sanctions — Nico's decision)
# --------------------------------------------------------------------------

def slide_10_nickel(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_OLIGO)
    _draw_action_title(slide, "Oligopoly Example: Nickel Mining")

    bullets = [
        ("Homogeneous good", 0),
        ("Key input for batteries (electric vehicles…)", 1),
        ("Two dominant producers – our duopoly example", 0),
        ("Vale (Brazil)", 1),
        ("Tsingshan Holding (China) – mines mainly in Indonesia", 1),
        ("A few other, smaller producers in different countries", 0),
        ("Will use this as guiding example", 0),
    ]
    _add_hierarchical_bullets(slide, left=MARGIN, top=Inches(1.85),
                              width=Inches(7.3), height=Inches(3.9),
                              items=bullets, size=24, sub_size=22,
                              line_spacing_pts=12)

    pic = _add_media_image(slide, "image21.jpeg", left=Inches(8.05),
                           top=Inches(2.35), width=Inches(4.7),
                           rounded=True, shadow=True)

    # note card narrowed + moved down by hand (Nico 2026-08-04; rendered
    # coords decoded from his grouped/scaled shapes)
    _add_convention_box(
        slide, MARGIN, Inches(6.554), Inches(10.965), Inches(0.470),
        prefix="Note:  ",
        body="Norilsk Nickel (Russia) is also a major producer – we "
             "exclude it here because of the sanctions against Russia",
        size=16, pad_h=Inches(0.12), pad_v=Inches(0.0))

    _draw_footer(slide, FOOTER_TEXT, 10)
    _set_notes(slide, (
        "Nickel is our guiding example for oligopoly: a homogeneous "
        "commodity, and a key input for nickel-rich EV batteries. We "
        "stylize the market as a duopoly with the two dominant producers: "
        "Vale of Brazil, the largest Western producer, and China's "
        "Tsingshan Holding, the world's largest nickel producer, which "
        "mines mainly in Indonesia — Indonesia now accounts for roughly "
        "two-thirds of world mine production (USGS 2026), and Tsingshan "
        "alone for roughly 30%. Norilsk Nickel of Russia is also a major "
        "producer, but we exclude it here because of the sanctions "
        "against Russia. Background on the nickel market: "
        "https://www.nsenergybusiness.com/features/"
        "top-nickel-producing-companies/"))
    return slide


# --------------------------------------------------------------------------
# Slide 20 — duopoly setup (Vale = Firm A, Tsingshan = Firm B)
# --------------------------------------------------------------------------

def slide_21_duopoly_setup(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_COURNOT)
    _draw_action_title(slide, "Duopoly Example: Nickel")

    eq_qm = (_oQ('M') + _omml_text(' = ') + _oQ('A')
             + _omml_text(' + ') + _oQ('B'))
    eq_dem = _omml_run('P') + _omml_text(' = 100 − ') + _oQ('M')

    segs = [
        ("text", "▪  Two firms: Vale (Firm A) and Tsingshan (Firm B)", {}),
        ("break", None, {}), ("break", None, {}),
        ("break", None, {}), ("break", None, {}),
        ("break", None, {}),
        ("text", "▪  Homogeneous (identical) good", {}),
        ("break", None, {}), ("break", None, {}),
        ("text", "▪  Market supply is the sum of the two firms' supply:",
         {}),
        ("break", None, {}),
        ("text", "        ", {}),
        ("omml", eq_qm, {'size': 26}),
        ("break", None, {}), ("break", None, {}),
        ("text", "▪  Market demand is:   ", {}),
        ("omml", eq_dem, {'size': 26}),
        ("break", None, {}), ("break", None, {}),
        ("text", "▪  The two firms have the same marginal cost of $10 "
                 "per kg", {}),
    ]
    # hand layout (Nico 2026-08-05): text block per his geometry,
    # logos side by side in the gap under the first bullet
    _add_mixed_textbox(slide, Inches(0.656), Inches(1.72),
                       Inches(9.804), Inches(4.982), segs,
                       default_size=24, default_color=NAVY)
    _add_media_image(slide, "image28.jpg", left=Inches(2.48),
                     top=Inches(2.243), width=Inches(2.4),
                     rounded=False, shadow=False)
    card = _add_outlined_box(slide, Inches(5.466), Inches(2.243),
                             Inches(2.4), Inches(1.05), "",
                             line=RULE, fill=WHITE, line_w=0.75,
                             rounded=True, shadow=False)
    tfc = card.text_frame
    par = tfc.paragraphs[0]
    par.alignment = PP_ALIGN.CENTER
    r1 = par.add_run()
    r1.text = "TSINGSHAN"
    r1.font.name = "Calibri"
    r1.font.size = Pt(22)
    r1.font.bold = True
    r1.font.color.rgb = RGBColor(0x00, 0x6B, 0x3F)
    par2 = tfc.add_paragraph()
    par2.alignment = PP_ALIGN.CENTER
    r2 = par2.add_run()
    r2.text = "HOLDING GROUP"
    r2.font.name = "Calibri"
    r2.font.size = Pt(11)
    r2.font.color.rgb = GRAY

    _draw_footer(slide, FOOTER_TEXT, 21)
    _set_notes(slide, (
        "The setup for our Cournot analysis. Two firms — Vale as Firm A "
        "and Tsingshan as Firm B — produce an identical good, nickel. "
        "Market supply is the sum of the two firms' quantities, market "
        "demand is P = 100 minus total quantity, and both firms have the "
        "same marginal cost of $10 per kg. All numbers are stylized to "
        "keep the algebra clean."))
    return slide


# --------------------------------------------------------------------------
# Slides 22/23 — Cournot reaction-function construction (parametrized)
# --------------------------------------------------------------------------

ACC6_75 = RGBColor(0xB9, 0x70, 0x34)   # Nico: accent6 lumMod 75%
ACC3_50 = RGBColor(0x4D, 0x5D, 0x2C)   # Nico: accent3 lumMod 50%
ACC1_BLUE = RGBColor(0x4F, 0x81, 0xBD)


def _cournot_reaction_slide(prs, page_num, *, roman, qb, firm_int,
                            opt_q, opt_p, mkt_q, lead_text, gap_y,
                            mkt_lbl, firm_lbl, firm_eq_full, mr_lbl,
                            units_pos, mask_firm_lbl=False,
                            pointer=None):
    """Residual-demand chart — ALL geometry/colors from Nico's
    hand-edits in backup_2026-08-05 (ported 2026-08-05, second pass)."""
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_COURNOT)
    _draw_action_title(slide,
                       f"Cournot Competition: Reaction Function ({roman})")

    _add_graphicframe_shadow(slide, Inches(0.42), Inches(1.5),
                             Inches(12.5), Inches(5.45))
    fig = SimpleFig(1.55, 6.35, 8.1, 4.45, 112, 112)
    _fig_axes(slide, fig)
    _add_text(slide, Inches(1.15), Inches(1.35), Inches(0.7),
              Inches(0.35), "P", size=18, bold=True, italic=True,
              color=NAVY, font="Calibri", align=PP_ALIGN.CENTER)
    _add_text(slide, Inches(9.7), Inches(6.41), Inches(0.6),
              Inches(0.35), "Q", size=18, bold=True, italic=True,
              color=NAVY, font="Calibri")

    # MC line (gray)
    _add_arrow(slide, (fig.x(0), fig.y(10)), (fig.x(105), fig.y(10)),
               color=GRAY, weight_pt=2.5, head=False)
    _add_text(slide, Inches(8.71), Inches(5.593), Inches(0.9),
              Inches(0.32), "MC", size=16, bold=True, color=GRAY,
              font="Calibri", align=PP_ALIGN.CENTER)

    # market demand (navy)
    _add_arrow(slide, (fig.x(0), fig.y(100)), (fig.x(100), fig.y(0)),
               color=NAVY, weight_pt=2.5, head=False)
    _add_mixed_textbox(
        slide, Inches(mkt_lbl[0]), Inches(mkt_lbl[1]), Inches(2.9),
        Inches(0.8),
        [("text", "Market Demand", {'bold': True, 'size': 17}),
         ("break", None, {}),
         ("omml", _omml_run('P') + _omml_text(' = 100 − ') + _oQ('M'),
          {'size': 16})],
        default_size=16, default_color=NAVY)

    # firm A residual demand + MR (red)
    _add_arrow(slide, (fig.x(0), fig.y(firm_int)),
               (fig.x(firm_int), fig.y(0)),
               color=FIRM_A_RED, weight_pt=2.5, head=False)
    firm_segs = [("text", "Firm A's Demand",
                  {'bold': True, 'size': 16, 'color': FIRM_A_RED})]
    if firm_eq_full:
        firm_segs += [
            ("break", None, {}),
            ("omml", _omml_run('P') + _omml_text(' = 100 − (')
             + _oQ('A') + _omml_text(f' + {qb})'), {'size': 14}),
            ("break", None, {}),
            ("omml", _omml_text(f'    = {firm_int} − ') + _oQ('A'),
             {'size': 14}),
        ]
    else:
        firm_segs += [
            ("break", None, {}),
            ("omml", _omml_run('P') + _omml_text(f' = {firm_int} − ')
             + _oQ('A'), {'size': 14}),
        ]
    firm_box = _add_mixed_textbox(
        slide, Inches(firm_lbl[0]), Inches(firm_lbl[1]),
        Inches(firm_lbl[2]),
        Inches(firm_lbl[3] if len(firm_lbl) > 3 else 1.1),
        firm_segs, default_size=14, default_color=FIRM_A_RED)
    if mask_firm_lbl:
        firm_box.fill.solid()
        firm_box.fill.fore_color.rgb = WHITE
    _add_arrow(slide, (fig.x(0), fig.y(firm_int)),
               (fig.x(firm_int / 2), fig.y(0)),
               color=FIRM_A_RED, weight_pt=2.0, head=False, dash='dash')
    _add_math_equation(slide, Inches(mr_lbl[0]), Inches(mr_lbl[1]),
                       Inches(0.95), Inches(0.36), _oMR('A'),
                       size_pt=15, color=FIRM_A_RED)
    if pointer is not None:
        (px1, py1), (px2, py2) = pointer
        _add_arrow(slide, (Inches(px1), Inches(py1)),
                   (Inches(px2), Inches(py2)), color=FIRM_A_RED,
                   weight_pt=1.0, head=True, head_size='sm')

    # "qb units" gap arrow between the demand curves (red, head left)
    _add_arrow(slide, (fig.x(100 - gap_y), fig.y(gap_y)),
               (fig.x(firm_int - gap_y), fig.y(gap_y)),
               color=FIRM_A_RED, weight_pt=1.75, head=True)
    _add_text(slide, Inches(units_pos[0]), Inches(units_pos[1]),
              Inches(1.6), Inches(0.3), f"{qb} units", size=14,
              bold=True, color=FIRM_A_RED, font="Calibri",
              align=PP_ALIGN.CENTER)

    # optimum: MR = MC at opt_q
    _add_arrow(slide, (fig.x(opt_q), fig.y(0)), (fig.x(opt_q),
               fig.y(opt_p)), color=GRAY, weight_pt=1.25, head=False,
               dash='dash')
    _add_arrow(slide, (fig.x(0), fig.y(opt_p)), (fig.x(mkt_q),
               fig.y(opt_p)), color=GRAY, weight_pt=1.25, head=False,
               dash='dash')
    _fig_point(slide, fig, opt_q, 10, fill=NAVY)
    _fig_point(slide, fig, opt_q, opt_p, fill=GOLD, line=NAVY)

    # ticks — optimal quantity/price highlighted per Nico's colors
    _fig_ytick(slide, fig, 100, "$100")
    _fig_ytick(slide, fig, firm_int, f"${firm_int}",
               color=FIRM_A_RED)
    _fig_ytick(slide, fig, opt_p, f"${opt_p}", bold=True)
    _fig_ytick(slide, fig, 10, "$10", color=GRAY)
    _fig_xtick(slide, fig, opt_q, str(opt_q), bold=True,
               color=FIRM_A_RED)
    for xv in sorted({firm_int, 100}):
        if abs(xv - opt_q) > 6:
            _fig_xtick(slide, fig, xv, str(xv))
    # market-quantity guide + dot + tick = ONE group with its OWN
    # final click (Nico's hand group, 2026-08-12)
    g1 = _add_arrow(slide, (fig.x(mkt_q), fig.y(opt_p)), (fig.x(mkt_q),
                    fig.y(0)), color=GRAY, weight_pt=1.25, head=False,
                    dash='dash')
    g2 = _fig_point(slide, fig, mkt_q, opt_p, fill=NAVY)
    g3 = _fig_xtick(slide, fig, mkt_q, str(mkt_q))
    _group_shapes(slide, [g1, g2, g3], "MktGuide")

    # assumption callout: white rounded card, red border, RED text with
    # a real Q_B subscript (box + text = one unit → group in phase 3)
    _add_outlined_box(slide, Inches(9.28), Inches(2.33), Inches(3.47),
                      Inches(1.0), "", line=FIRM_A_RED, fill=WHITE,
                      line_w=1.25, rounded=True, shadow=False,
                      corner_pct=0.16)
    ctb = slide.shapes.add_textbox(Inches(9.44), Inches(2.477),
                                   Inches(3.15), Inches(0.71))
    ctf = ctb.text_frame
    ctf.word_wrap = True
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    cpara = ctf.paragraphs[0]

    def _crun(t, bold=False, sub=False, color=FIRM_A_RED):
        r = cpara.add_run()
        r.text = t
        r.font.name = "Calibri"
        r.font.size = Pt(18)
        r.font.bold = bold
        r.font.color.rgb = color
        if sub:
            r._r.get_or_add_rPr().set('baseline', '-25000')

    _crun(lead_text + " (")
    _crun("Q", bold=True, color=ACC3_50)
    _crun("B", bold=True, sub=True, color=ACC3_50)
    _crun(f" = {qb}", bold=True, color=ACC3_50)
    _crun(")")

    _add_text(slide, Inches(9.95), Inches(3.45), Inches(2.33),
              Inches(0.9),
              f"→  Firm A's optimal response is {opt_q} units",
              size=17, bold=True, color=FIRM_A_RED, font="Calibri")

    _draw_footer(slide, FOOTER_TEXT, page_num)
    _set_notes(slide, (
        "Cournot assumed quantity competition: each firm sets quantity "
        "taking as fixed the quantity set by the other. The assumption is "
        "that the other firm will continue to produce the same quantity "
        "whatever I do. [Of course most of the time this assumption will "
        "prove false, but not at the equilibrium]. The residual demand "
        "curve is just the amount of demand left over when subtracting "
        "the other firm's fixed supply from market demand. As we will "
        "see shortly price does not go to the competitive level but "
        "settles between the monopoly and perfect competition price.\n"
        f"Step 1: Go from the market-level demand to the individual-level "
        f"demand by 'assuming' QB={qb}: P=100−(QA+{qb})={firm_int}−QA\n"
        f"Step 2: Once you have the individual-level demand, find the "
        f"profit-maximizing quantity as usual (MR=MC): "
        f"MRA={firm_int}−2QA; {firm_int}−2QA=10 → QA={opt_q}."))
    return slide

def slide_23_reaction_i(prs):
    return _cournot_reaction_slide(
        prs, 23, roman="I", qb=50, firm_int=50, opt_q=20, opt_p=30,
        mkt_q=70, gap_y=24, units_pos=(4.577, 5.494),
        mkt_lbl=(3.358, 3.0), firm_lbl=(1.804, 4.089, 3.0),
        firm_eq_full=True, mr_lbl=(2.996, 6.427),
        lead_text="Firm A assumes that B produces 50 units")

def slide_24_reaction_ii(prs):
    return _cournot_reaction_slide(
        prs, 24, roman="II", qb=20, firm_int=80, opt_q=35, opt_p=45,
        mkt_q=55, gap_y=35, units_pos=(4.842, 4.989),
        mkt_lbl=(6.401, 4.563), firm_lbl=(2.342, 2.793, 2.46, 0.50),   # small box, Nico 2026-08-12
        firm_eq_full=False, mr_lbl=(4.135, 6.438),
        mask_firm_lbl=True,
        pointer=((2.566, 3.171), (2.273, 3.565)),
        lead_text="Now Firm A assumes that B produces 20 units")

# --------------------------------------------------------------------------

def slide_25_reaction_function(prs):
    """All geometry/colors from Nico's hand-edits (backup_2026-08-05):
    orange (accent6@75%) reaction points, guides, tick labels and RP
    labels with pointer ticks; olive (accent3@50%) x-axis title."""
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_COURNOT)
    _draw_action_title(slide, "Firm A's Reaction Function")

    _add_graphicframe_shadow(slide, Inches(1.7), Inches(1.55),
                             Inches(10.037), Inches(5.04))
    fig = SimpleFig(3.35, 5.8, 7.0, 3.75, 100, 55)
    _fig_axes(slide, fig)
    ttl = _add_mixed_textbox(
        slide, Inches(1.033), Inches(3.315), Inches(3.4), Inches(0.32),
        [("text", "Firm A's Reaction (", {'bold': True, 'italic': True,
                                          'size': 16,
                                          'color': FIRM_A_RED}),
         ("omml", _oQ('A'), {'size': 16, 'color': FIRM_A_RED}),
         ("text", ")", {'bold': True, 'italic': True, 'size': 16,
                        'color': FIRM_A_RED})],
        default_size=16, default_color=FIRM_A_RED)
    ttl.rotation = 270
    _add_mixed_textbox(
        slide, Inches(5.45), Inches(6.16), Inches(3.2), Inches(0.32),
        [("text", "Firm B's Output (", {'bold': True, 'italic': True,
                                        'size': 16,
                                        'color': ACC3_50}),
         ("omml", _oQ('B'), {'size': 16, 'color': ACC3_50}),
         ("text", ")", {'bold': True, 'italic': True, 'size': 16,
                        'color': ACC3_50})],
        default_size=16, default_color=ACC3_50)

    # reaction function QA = 45 − 0.5 QB (red)
    _add_arrow(slide, (fig.x(0), fig.y(45)), (fig.x(90), fig.y(0)),
               color=FIRM_A_RED, weight_pt=2.5, head=False)
    _add_mixed_textbox(
        slide, Inches(7.62), Inches(4.525), Inches(3.3), Inches(0.6),
        [("text", "Firm A's reaction function", {'bold': True,
                                                 'size': 16,
                                                 'color': FIRM_A_RED}),
         ("break", None, {}),
         ("omml", _omml_sub(_omml_run('Q', color=FIRM_A_RED),
                            _omml_text('A', color=FIRM_A_RED))
          + _omml_text(' = 45 − 0.5·', color=FIRM_A_RED)
          + _omml_sub(_omml_run('Q', color=ACC3_50),
                      _omml_text('B', color=ACC3_50)),
          {'size': 16})],
        default_size=16, default_color=FIRM_A_RED)

    # reaction points: orange dotted guides, gold dots, orange ticks
    for (qb, qa) in [(50, 20), (20, 35)]:
        _add_arrow(slide, (fig.x(qb), fig.y(0)), (fig.x(qb), fig.y(qa)),
                   color=ACC6_75, weight_pt=1.25, head=False,
                   dash='sysDot')
        _add_arrow(slide, (fig.x(qb), fig.y(qa)), (fig.x(0), fig.y(qa)),
                   color=ACC6_75, weight_pt=1.25, head=False,
                   dash='sysDot')
        _fig_point(slide, fig, qb, qa, fill=GOLD, line=NAVY,
                   r_in=0.065)
        _fig_xtick(slide, fig, qb, str(qb), bold=True, color=ACC6_75)
        _fig_ytick(slide, fig, qa, str(qa), bold=True, color=ACC6_75)
    _add_text(slide, Inches(4.466), Inches(2.792), Inches(2.2),
              Inches(0.28), "Reaction point (II)", size=14, bold=True,
              italic=True, color=ACC6_75, font="Calibri")
    _add_arrow(slide, (Inches(4.919), Inches(3.062)),
               (Inches(4.807), Inches(3.372)), color=ACC6_75,
               weight_pt=1.0, head=True, head_size='sm')
    _add_text(slide, Inches(6.383), Inches(3.782), Inches(2.2),
              Inches(0.28), "Reaction point (I)", size=14, bold=True,
              italic=True, color=ACC6_75, font="Calibri")
    _add_arrow(slide, (Inches(6.971), Inches(4.062)),
               (Inches(6.859), Inches(4.372)), color=ACC6_75,
               weight_pt=1.0, head=True, head_size='sm')
    _fig_ytick(slide, fig, 45, "45")
    _fig_xtick(slide, fig, 90, "90")

    _draw_footer(slide, FOOTER_TEXT, 25)
    _add_video_link_box(slide,
                        "Practice Video “Cournot Competition Math”",
                        size=16)
    _set_notes(slide, (
        "Don't say they move sequentially — it's just 'contemplating'. "
        "The reaction function collects Firm A's optimal responses: for "
        "every quantity B might produce, it gives A's profit-maximizing "
        "quantity. The two points we just derived — B produces 50, A "
        "responds with 20; B produces 20, A responds with 35 — pin down "
        "the line QA = 45 − 0.5·QB. The same logic (with the general "
        "algebra) is worked through in the practice video."))
    return slide

# --------------------------------------------------------------------------

def slide_26_cournot_equilibrium(prs):
    """All geometry/colors from Nico's hand-edits (backup_2026-08-05):
    olive Firm-B labels + 2.5 pt olive pointer, black A-equation,
    subscripted 'Starting point' text; equilibrium callout in the gold
    action-box format (newer request, kept)."""
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_COURNOT)
    _draw_action_title(slide, "Cournot Equilibrium")

    _add_graphicframe_shadow(slide, Inches(0.45), Inches(1.5),
                             Inches(12.65), Inches(5.45))
    fig = SimpleFig(2.85, 6.3, 6.7, 4.4, 100, 100)
    _fig_axes(slide, fig)
    ttl = _add_mixed_textbox(
        slide, Inches(0.728), Inches(3.795), Inches(3.0), Inches(0.32),
        [("text", "Firm A's Output (", {'bold': True, 'italic': True,
                                        'size': 16,
                                        'color': FIRM_A_RED}),
         ("omml", _oQ('A'), {'size': 16, 'color': FIRM_A_RED}),
         ("text", ")", {'bold': True, 'italic': True, 'size': 16,
                        'color': FIRM_A_RED})],
        default_size=16, default_color=FIRM_A_RED)
    ttl.rotation = 270
    _add_mixed_textbox(
        slide, Inches(4.7), Inches(6.599), Inches(3.0), Inches(0.32),
        [("text", "Firm B's Output (", {'bold': True, 'italic': True,
                                        'size': 16,
                                        'color': ACC3_50}),
         ("omml", _oQ('B'), {'size': 16, 'color': ACC3_50}),
         ("text", ")", {'bold': True, 'italic': True, 'size': 16,
                        'color': ACC3_50})],
        default_size=16, default_color=ACC3_50)

    # Firm A's reaction function (red line, BLACK equation per Nico)
    _g = _add_arrow(slide, (fig.x(0), fig.y(45)), (fig.x(90), fig.y(0)),
               color=FIRM_A_RED, weight_pt=2.5, head=False)
    _gt = _add_mixed_textbox(
        slide, Inches(6.334), Inches(4.98), Inches(3.0), Inches(0.85),
        [("text", "Firm A's reaction function", {'bold': True,
                                                 'size': 14,
                                                 'color': FIRM_A_RED}),
         ("break", None, {}),
         ("omml", _omml_sub(_omml_run('Q', color=FIRM_A_RED),
                            _omml_text('A', color=FIRM_A_RED))
          + _omml_text(' = 45 − 0.5·', color=FIRM_A_RED)
          + _omml_sub(_omml_run('Q', color=ACC3_50),
                      _omml_text('B', color=ACC3_50)),
          {'size': 14})],
        default_size=14, default_color=FIRM_A_RED)
    _group_shapes(slide, [_g, _gt], "ReactionA")
    # Firm B's reaction function (green line, olive label)
    _g = _add_arrow(slide, (fig.x(0), fig.y(90)), (fig.x(45), fig.y(0)),
               color=FIRM_B_GREEN, weight_pt=2.5, head=False)
    _gt = _add_mixed_textbox(
        slide, Inches(3.365), Inches(2.301), Inches(3.0), Inches(0.6),
        [("text", "Firm B's reaction function", {'bold': True,
                                                 'size': 14,
                                                 'color': ACC3_50}),
         ("break", None, {}),
         ("omml", _omml_sub(_omml_run('Q', color=ACC3_50),
                            _omml_text('B', color=ACC3_50))
          + _omml_text(' = 45 − 0.5·', color=ACC3_50)
          + _omml_sub(_omml_run('Q', color=FIRM_A_RED),
                      _omml_text('A', color=FIRM_A_RED)),
          {'size': 14})],
        default_size=14, default_color=ACC3_50)
    # olive pointer arrow from B's label to its line (2.5 pt, open head)
    ptr = _add_arrow(slide, (Inches(3.842), Inches(2.689)),
                     (Inches(3.343), Inches(2.915)), color=ACC3_50,
                     weight_pt=2.5, head=True)
    ln = ptr.line._get_or_add_ln()
    te = ln.find(qn('a:tailEnd'))
    te.set('type', 'arrow')
    _group_shapes(slide, [_g, _gt, ptr], "ReactionB")

    # convergence staircase from QA = 15 (tick + first segment = one
    # group per Nico 2026-08-12; the other segments click one by one)
    _t15 = _fig_ytick(slide, fig, 15, "15")
    path = [(0, 15), (37.5, 15), (37.5, 26.25), (31.9, 26.25),
            (31.9, 29.1)]
    _segs = []
    for (x1, y1), (x2, y2) in zip(path, path[1:]):
        _segs.append(_add_arrow(slide, (fig.x(x1), fig.y(y1)),
                     (fig.x(x2), fig.y(y2)), color=GRAY,
                     weight_pt=1.5, head=True, dash='sysDash'))
    _group_shapes(slide, [_t15, _segs[0]], "Fifteen")

    # equilibrium guides + the two "30" ticks = one group
    _ga = _add_arrow(slide, (fig.x(30), fig.y(0)), (fig.x(30),
               fig.y(30)), color=GRAY, weight_pt=1.25, head=False,
               dash='dash')
    _gb = _add_arrow(slide, (fig.x(0), fig.y(30)), (fig.x(30),
               fig.y(30)), color=GRAY, weight_pt=1.25, head=False,
               dash='dash')
    _gc = _fig_xtick(slide, fig, 30, "30", bold=True)
    _gd = _fig_ytick(slide, fig, 30, "30", bold=True)
    _group_shapes(slide, [_ga, _gb, _gc, _gd], "Thirty")
    _fig_xtick(slide, fig, 45, "45")
    _fig_xtick(slide, fig, 90, "90")
    _fig_ytick(slide, fig, 45, "45")
    _fig_ytick(slide, fig, 90, "90")

    # equilibrium point + arrow + gold box + OMML = one group
    _e1 = _fig_point(slide, fig, 30, 30, fill=GOLD, line=NAVY,
                     r_in=0.075)
    _e2 = _add_arrow(slide, (Inches(5.195), Inches(4.32)),
               (Inches(4.927), Inches(4.892)), color=NAVY,
               weight_pt=1.5, head=True)
    _e3 = _add_rounded_filled_box(
        slide, Inches(5.16), Inches(3.78), Inches(3.12), Inches(0.73),
        "", fill=GOLD, text_color=NAVY, size=14, bold=True,
        corner_pct=0.12, shadow=True)
    _e4 = _add_mixed_textbox(
        slide, Inches(5.262), Inches(3.836), Inches(3.1), Inches(0.65),
        [("text", "Cournot Equilibrium", {'bold': True, 'size': 15}),
         ("break", None, {}),
         ("omml", _omml_sub(_omml_run('Q', color=FIRM_A_RED),
                            _omml_text('A', color=FIRM_A_RED))
          + _omml_text(' = ')
          + _omml_sub(_omml_run('Q', color=ACC3_50),
                      _omml_text('B', color=ACC3_50))
          + _omml_text(' = 30  and  ') + _omml_run('P')
          + _omml_text(' = 40'), {'size': 14,
                                  'color': RGBColor(0, 0, 0)})],
        default_size=14, default_color=NAVY)
    _group_shapes(slide, [_e1, _e2, _e3, _e4], "Equilibrium")

    # callout boxes (top-right)
    _s26_a = _shape_ids(slide)
    _add_convention_box(
        slide, Inches(8.35), Inches(1.7), Inches(4.6), Inches(1.25),
        runs=[("Cournot equilibrium: neither firm would like to "
               "deviate", {'size': 16, 'color': NAVY}),
              ("→  Look for intersection of reaction functions",
               {'size': 16, 'bold': True, 'color': NAVY,
                'newline': True})],
        fill_rgb=WHITE, border=NAVY, line_w=1.25, size=16)
    _group_shapes(slide, _new_shapes_since(slide, _s26_a), "Callout")
    _s26_b = _shape_ids(slide)
    _add_convention_box(
        slide, Inches(8.85), Inches(3.2), Inches(4.1), Inches(0.95),
        runs=[("STARTPLACEHOLDER", {'size': 15, 'italic': True,
                                    'color': GRAY})],
        fill_rgb=WHITE, border=GRAY, line_w=1.0, size=15)
    # rewrite the starting-point text with a real Q_A subscript
    for sh in slide.shapes:
        if sh.has_text_frame and "STARTPLACEHOLDER" in sh.text_frame.text:
            para = sh.text_frame.paragraphs[0]
            for r in list(para.runs):
                r._r.getparent().remove(r._r)

            def _srun(t, sub=False):
                r = para.add_run()
                r.text = t
                r.font.name = "Calibri"
                r.font.size = Pt(15)
                r.font.italic = True
                r.font.color.rgb = GRAY
                if sub:
                    r._r.get_or_add_rPr().set('baseline', '-25000')

            _srun("Starting point: Firm B contemplates its reaction "
                  "to Q")
            _srun("A", sub=True)
            _srun(" = 15")
    _group_shapes(slide, _new_shapes_since(slide, _s26_b), "Starting")

    _draw_footer(slide, FOOTER_TEXT, 26)
    _set_notes(slide, (
        "Don't say they move sequentially — it's just 'contemplating'. "
        "The last step is to figure out the Cournot equilibrium, given by "
        "the 'intersection' of the two reaction functions. The intuition: "
        "we want a pair (QA*, QB*) so that A wants to respond to QB* with "
        "QA*, and B wants to respond to QA* by picking QB*.\n"
        "We start from the reaction function for A: QA=45−0.5QB (1). "
        "Then we repeat the analysis for firm B (everything is symmetric, "
        "so it is the mirror image): QB=45−0.5QA (2). Last, we need the "
        "pair (QA*, QB*) that satisfies both — two equations, two "
        "unknowns. Replace (1) in (2): QB=45−0.5(45−0.5QB)=22.5+0.25QB → "
        "0.75QB=22.5 → QB=30. By symmetry QA=30 as well (verify: "
        "QA=45−0.5·30=30). The figure shows the graphical version: two "
        "lines with different slopes intersect in exactly one point. The "
        "staircase from QA=15 shows how the contemplated responses "
        "converge to the equilibrium."))
    return slide

# --------------------------------------------------------------------------

def slide_27_computation(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_COURNOT)
    _draw_action_title(slide, "Cournot: Math")   # renamed by Nico 2026-08-13

    box_w = Inches(7.2)
    left = (SLIDE_W - box_w) // 2
    _draw_footer(slide, FOOTER_TEXT, 27)
    _add_video_link_box(slide,
                        "Practice Video “Cournot Competition Math”",
                        left=left, top=Inches(2.5), width=box_w,
                        height=Inches(0.85), size=20)
    _add_text(slide, MARGIN, Inches(4.265), RULE_W, Inches(0.4),
              "More general case with different MC:",
              size=24, bold=False, color=NAVY, font="Calibri",
              align=PP_ALIGN.CENTER)
    _add_video_link_box(slide,
                        "Practice Video “Oligopoly with different MC”",
                        left=left, top=Inches(4.75), width=box_w,
                        height=Inches(0.85), size=20)
    _set_notes(slide, (
        "Two practice videos on BruinLearn walk through the Cournot "
        "algebra step by step: the baseline case we just did, and the "
        "more general case where the two firms have different marginal "
        "costs."))
    return slide

def slide_28_further_examples(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_COURNOT)
    _draw_action_title(
        slide, "Further Examples for Cournot Competition")

    bullets = [
        ("Oil", 0),
        ("However: Includes a cartel (OPEC)", 1),
        ("Steel producers", 0),
    ]
    _add_hierarchical_bullets(slide, left=MARGIN, top=Inches(2.5),
                              width=Inches(6.6), height=Inches(3.0),
                              items=bullets, size=26, sub_size=24,
                              line_spacing_pts=20)

    _add_media_image(slide, "image30.jpg", left=Inches(6.09),
                     top=Inches(1.41), width=Inches(3.98),
                     rounded=False, shadow=True)
    # native pie chart; layout hand-tuned by Nico 2026-08-05
    _add_graphicframe_shadow(slide, Inches(3.39), Inches(3.83),
                             Inches(6.0), Inches(3.05))
    _add_text(slide, Inches(3.39), Inches(3.94), Inches(6.0),
              Inches(0.3),
              "US crude steel production by producer, 2024",
              size=13, italic=True, bold=True, color=NAVY,
              font="Calibri", align=PP_ALIGN.CENTER)
    chart_data = CategoryChartData()
    chart_data.categories = ["Nucor", "Cleveland-Cliffs",
                             "U.S. Steel (Nippon Steel)",
                             "Steel Dynamics", "CMC", "Others"]
    chart_data.add_series('Share', (26, 20, 13, 13, 6, 22))
    gframe = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, Inches(3.54), Inches(4.22), Inches(5.7),
        Inches(2.58), chart_data)
    chart = gframe.chart
    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(10)
    chart.legend.font.name = "Carlito"
    chart.legend.font.color.rgb = NAVY
    plot = chart.plots[0]
    plot.has_data_labels = True
    dlab = plot.data_labels
    dlab.show_percentage = True
    dlab.show_value = False
    dlab.show_category_name = False
    dlab.font.size = Pt(10)
    dlab.font.bold = True
    dlab.font.name = "Carlito"
    dlab.font.color.rgb = WHITE
    slice_colors = [NAVY, GOLD, RGBColor(0xC0, 0x50, 0x4D),
                    RGBColor(0x1B, 0x5E, 0x20), RGBColor(0x55, 0x5B,
                    0x66), RGBColor(0xC8, 0xCD, 0xD3)]
    for pt_i, col in zip(plot.series[0].points, slice_colors):
        pt_i.format.fill.solid()
        pt_i.format.fill.fore_color.rgb = col
    _add_text(slide, Inches(3.39), Inches(6.61), Inches(6.0),
              Inches(0.25),
              "Source: worldsteel, company reports. U.S. Steel acquired "
              "by Nippon Steel, June 2025",
              size=9, italic=True, color=GRAY, font="Calibri",
              align=PP_ALIGN.CENTER)

    _draw_footer(slide, FOOTER_TEXT, 28)
    _set_notes(slide, (
        "Why OPEC is not sued for antitrust: '[i]f the members of OPEC "
        "were private companies and not nations, they long ago would have "
        "been prosecuted for engaging in illegal price fixing.' However, "
        "OPEC member nations have argued that doctrines of sovereign "
        "immunity and act of state make them immune from U.S. antitrust "
        "prosecution for their otherwise illegal price fixing cartel. "
        "Indeed, a federal court decision in 1979 ruled that the actions "
        "of OPEC were immune from U.S. antitrust scrutiny because of the "
        "doctrine of sovereign immunity (the proposition that the "
        "government cannot be sued without its consent).\n"
        "Source: https://energyfuse.org/should-we-still-be-concerned-"
        "about-opec/\n"
        "Steel pie (updated 2026): shares of US crude steel production "
        "2024, total 79.5 Mt (worldsteel). Nucor 26% (20.7 Mt), "
        "Cleveland-Cliffs 20% (~16 Mt - it absorbed BOTH ArcelorMittal "
        "USA and AK Steel in 2020, which is why the old chart looks so "
        "different), U.S. Steel 13% (10.7 Mt, acquired by Nippon Steel "
        "June 2025, still operating as U.S. Steel), Steel Dynamics 13% "
        "(10.2 Mt), CMC ~6% (approx., fiscal-year basis), Others ~22% "
        "(Gerdau, North Star BlueScope, NLMK USA, SSAB Americas, etc). "
        "Sources: worldsteel company rankings 2024; U.S. Steel 10-K "
        "segment data; company reports."))
    return slide


# --------------------------------------------------------------------------
# Slide 62 — trade wars (UPDATED 2026-07-29: US vs. China labels; original
# stylized payoffs kept; calibration details in the speaker notes)
# --------------------------------------------------------------------------

def slide_63_trade_wars(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_GAMES)
    _draw_action_title(slide, "Game Theory and Trade Wars")

    _add_text(slide, MARGIN + Inches(0.1), Inches(1.62), Inches(12.4),
              Inches(0.45),
              "In a trade war, each country's dominant strategy is to set "
              "high tariffs",
              size=24, color=NAVY, font="Calibri")

    anchors = _add_payoff_matrix(
        slide, left=Inches(4.35), top=Inches(3.05),
        cell_w=Inches(2.35), cell_h=Inches(1.05),
        row_player="US", col_player="China",
        row_strats=["Low Tariffs", "High Tariffs"],
        col_strats=["Low Tariffs", "High Tariffs"],
        payoffs=[[("20", "20"), ("13", "24")],
                 [("24", "13"), ("15", "15")]],
        caption=True, nash_cells=[(1, 1)])

    # Nash-equilibrium callout with connector to the equilibrium cell
    _add_outlined_box(
        slide, Inches(10.15), Inches(3.75), Inches(2.65), Inches(0.62),
        "Nash Equilibrium", line=GOLD, text_color=NAVY, size=18,
        bold=True, line_w=1.75, rounded=True, shadow=True,
        corner_pct=0.25)
    cx, cy = anchors[(1, 1, 'cell')]
    _add_arrow(slide, (Inches(10.55), Inches(4.37)),
               (cx + Inches(1.05), cy - Inches(0.28)),
               color=GOLD, weight_pt=1.75, head=True)

    _add_text(slide, MARGIN + Inches(0.1), Inches(5.85), Inches(6.4),
              Inches(0.4), "Interpretation of payoffs:",
              size=20, underline=True, color=NAVY, font="Calibri")
    _add_text(slide, MARGIN + Inches(0.1), Inches(6.22), Inches(6.4),
              Inches(0.4), "Countries' real GDP ($ trillion, "
              "price-adjusted)", size=20, color=NAVY, font="Calibri")
    _add_text(slide, Inches(7.1), Inches(6.32), Inches(5.9), Inches(0.6),
              "2025–26 tariff war: US average tariff on Chinese goods "
              "≈ 48%; China's on US goods ≈ 32% (PIIE)",
              size=14, italic=True, color=GRAY, font="Calibri",
              align=PP_ALIGN.RIGHT)

    _draw_footer(slide, FOOTER_TEXT, 63)
    _set_notes(slide, (
        "The payoff numbers are stylized, but they are calibrated to real "
        "magnitudes. Interpretation of payoffs: real GDP in $ trillion. "
        "Actual 2025 nominal GDP: US ≈ $30.6 trillion, China ≈ $19.4 "
        "trillion (IMF World Economic Outlook, Oct 2025).\n"
        "Where we actually are: after the 2025 tariff escalations and the "
        "partial truces (Geneva, May 2025; Korea, Nov 2025), the average "
        "US tariff on Chinese imports is ≈ 47.5% and China's average "
        "tariff on US imports ≈ 31.9% (PIIE US–China trade war tariff "
        "tracker, Nov 2025) — i.e., the high-tariff / high-tariff cell.\n"
        "Estimated costs behind the payoff gaps: the 2025 tariff actions "
        "are estimated to reduce long-run US GDP by roughly 0.3–0.7% "
        "(Yale Budget Lab; Tax Foundation), on the order of $100–200B "
        "per year; for China, estimates run from 0.5 to 2.5 percentage "
        "points off GDP growth over 2025–27 depending on escalation "
        "(EIU).\n"
        "The game's structure matches the evidence: a unilateral tariff "
        "shifts some surplus to the imposing country (the 24 vs. 13 "
        "cells), but when both retaliate, both end up worse off (15, 15) "
        "than under mutual low tariffs (20, 20) — the prisoner's-dilemma "
        "logic."))
    return slide



# --------------------------------------------------------------------------
# Batch A/B helpers: editable Bézier freeform (custGeom), OMML tick labels,
# best-response arrows, poll badge shortcut, firm colors
# --------------------------------------------------------------------------

P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
FIRM_A_RED = RGBColor(0xC0, 0x00, 0x00)
FIRM_B_GREEN = RGBColor(0x1B, 0x5E, 0x20)
GREEN_FILL = RGBColor(0xC6, 0xE0, 0xB4)
RED_FILL = RGBColor(0xF2, 0xB8, 0xB8)
PALE_GOLD = RGBColor(0xF6, 0xE8, 0xC9)

_BEZ_ID = [9000]


def _add_bezier_curve(slide, start, segments, *, color=NAVY,
                      weight_pt=2.5, dash=None):
    """Open cubic-Bézier freeform with FEW anchors (Edit-Points friendly).
    start = (x, y); segments = [((c1x,c1y),(c2x,c2y),(px,py)), ...], all
    EMU ints/Inches."""
    pts = [start] + [p for seg in segments for p in seg]
    xs = [int(p[0]) for p in pts]
    ys = [int(p[1]) for p in pts]
    x0, y0 = min(xs), min(ys)
    w = max(max(xs) - x0, 1)
    h = max(max(ys) - y0, 1)

    def rel(p):
        return int(p[0]) - x0, int(p[1]) - y0

    sx, sy = rel(start)
    path = [f'<a:moveTo><a:pt x="{sx}" y="{sy}"/></a:moveTo>']
    for seg in segments:
        inner = "".join(f'<a:pt x="{rel(p)[0]}" y="{rel(p)[1]}"/>'
                        for p in seg)
        path.append(f'<a:cubicBezTo>{inner}</a:cubicBezTo>')
    dash_xml = f'<a:prstDash val="{dash}"/>' if dash else ''
    _BEZ_ID[0] += 1
    sp_xml = (
        f'<p:sp xmlns:p="{P_NS}" xmlns:a="{A_NS}">'
        f'<p:nvSpPr><p:cNvPr id="{_BEZ_ID[0]}" name="BezierCurve'
        f'{_BEZ_ID[0]}"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr><a:xfrm><a:off x="{x0}" y="{y0}"/>'
        f'<a:ext cx="{w}" cy="{h}"/></a:xfrm>'
        f'<a:custGeom><a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/>'
        f'<a:rect l="0" t="0" r="{w}" b="{h}"/>'
        f'<a:pathLst><a:path w="{w}" h="{h}" fill="none">'
        f'{"".join(path)}</a:path></a:pathLst></a:custGeom>'
        f'<a:noFill/>'
        f'<a:ln w="{int(weight_pt * 12700)}" cap="rnd">'
        f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
        f'{dash_xml}<a:round/></a:ln>'
        f'</p:spPr><p:txBody><a:bodyPr/><a:lstStyle/>'
        f'<a:p><a:endParaRPr lang="en-US"/></a:p></p:txBody></p:sp>')
    el = ET.fromstring(sp_xml)
    slide.shapes._spTree.append(el)
    return el


def _oq(s):
    return _omml_sub(_omml_run('q'), _omml_text(s))


def _fig_math_tick(slide, fig, axis, val, omml, *, color=NAVY,
                   size_pt=13):
    """Small OMML tick label (for subscripted P1/PC/Q1/QC/q1... ticks)."""
    if axis == 'y':
        left = Inches(fig.l) - Inches(0.72)
        top = fig.y(val) - Inches(0.16)
    else:
        left = fig.x(val) - Inches(0.30)
        top = Inches(fig.b + 0.05)
    _add_math_equation(slide, left, top, Inches(0.62), Inches(0.32),
                       omml, size_pt=size_pt, color=color)


def _add_cover_box(slide, left, top, width, height):
    """Light-blue gradient 'hidden answer' cover — the ORIGINAL deck's
    group-work / reveal technique (Nico 2026-08-09)."""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(left),
                                 int(top), int(width), int(height))
    shp.line.fill.background()
    shp.shadow.inherit = False
    shp.fill.gradient()
    stops = shp.fill.gradient_stops
    stops[0].color.rgb = RGBColor(0xF2, 0xF7, 0xFC)
    stops[0].position = 0.0
    stops[1].color.rgb = RGBColor(0xC9, 0xDC, 0xEE)
    stops[1].position = 1.0
    shp.fill.gradient_angle = 90
    return shp


def _br_arrow(slide, frm, to, color, *, weight_pt=2.5):
    """Best-response arrow between two payoff anchors, positioned like
    the ORIGINAL deck (Nico 2026-08-09): a vertical arrow runs
    number-to-number (starting just below/above the source payoff and
    stopping clear of the target's circle); a horizontal arrow sits
    ABOVE the numbers (e.g. from -1 to 0) so it never strikes through
    them."""
    (x1, y1), (x2, y2) = frm, to
    if x1 == x2:                        # vertical (row player)
        d = 1 if y2 > y1 else -1
        return _add_arrow(slide, (x1, y1 + d * Inches(0.26)),
                          (x2, y2 - d * Inches(0.40)), color=color,
                          weight_pt=weight_pt, head=True)
    d = 1 if x2 > x1 else -1            # horizontal (column player)
    yy = y1 - Inches(0.30)
    return _add_arrow(slide, (x1 + d * Inches(0.05), yy),
                      (x2 - d * Inches(0.36), yy), color=color,
                      weight_pt=weight_pt, head=True)


def _poll_badge(slide):
    _add_discussion_break(slide, text="Poll Break", width=Inches(3.0))


def _add_video_link_box(slide, text=None, *, left=None, top=None,
                        width=None, height=None, size=17):
    """Deck-standard Practice-Video link box (v2, Nico 2026-08-05):
    rounded card with the metallic VERTICAL GRAY GRADIENT (bg1 lumMod
    65% -> 95% -> 65%), gold 1.75 pt border, soft drop shadow, gold play
    glyph + navy bold label. DEFAULT position = bottom-right corner
    (6.92, 6.83, 5.85 x 0.58in — slide-24 reference position, overlaying
    the footer). Call AFTER _draw_footer so the box sits on top."""
    left = Inches(6.92) if left is None else left
    top = Inches(6.83) if top is None else top
    width = Inches(5.85) if width is None else width
    height = Inches(0.58) if height is None else height
    shp = _add_outlined_box(slide, left, top, width, height, "",
                            line=GOLD, fill=WHITE, line_w=1.75,
                            rounded=True, shadow=True, corner_pct=0.28)
    # swap the solid fill for Nico's vertical gray gradient
    spPr = shp._element.spPr
    for old in spPr.findall(qn('a:solidFill')):
        spPr.remove(old)
    grad = ET.fromstring(
        '<a:gradFill xmlns:a="http://schemas.openxmlformats.org/'
        'drawingml/2006/main"><a:gsLst>'
        '<a:gs pos="0"><a:schemeClr val="bg1">'
        '<a:lumMod val="65000"/></a:schemeClr></a:gs>'
        '<a:gs pos="54000"><a:schemeClr val="bg1">'
        '<a:lumMod val="95000"/></a:schemeClr></a:gs>'
        '<a:gs pos="100000"><a:schemeClr val="bg1">'
        '<a:lumMod val="65000"/></a:schemeClr></a:gs>'
        '</a:gsLst><a:lin ang="5400000" scaled="1"/></a:gradFill>')
    prstGeom = spPr.find(qn('a:prstGeom'))
    prstGeom.addnext(grad)
    par = shp.text_frame.paragraphs[0]
    r1 = par.runs[0] if par.runs else par.add_run()
    r1.text = "\u25b6  "
    r1.font.name = "Calibri"
    r1.font.size = Pt(size)
    r1.font.bold = True
    r1.font.color.rgb = GOLD
    r2 = par.add_run()
    r2.text = text
    r2.font.name = "Calibri"
    r2.font.size = Pt(size)
    r2.font.bold = True
    r2.font.color.rgb = NAVY
    return shp


def _dash_shape_line(shape, dash='dash'):
    ln = shape.line._get_or_add_ln()
    for old in ln.findall(qn('a:prstDash')):
        ln.remove(old)
    prst = ln.makeelement(qn('a:prstDash'), {'val': dash})
    fill = ln.find(qn('a:solidFill'))
    if fill is not None:
        fill.addnext(prst)
    else:
        ln.append(prst)
    return shape



# --------------------------------------------------------------------------
# Slide 12 — Cartels (two-panel diagram, native)
# --------------------------------------------------------------------------

def slide_13_cartels(prs):
    """Two-panel cartel diagram. Geometry, tick names and wording =
    Nico's hand-edits (2026-08-04, extended 2026-08-12). Topology
    mirrors his hand groupings: each price line / quantity guide /
    profit region travels with its tick(s) and label as ONE group."""
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_COLLU)
    _draw_action_title(slide, "Cartels")

    _add_graphicframe_shadow(slide, Inches(0.45), Inches(1.43),
                             Inches(12.45), Inches(4.45))

    def _tick(x_in, y_in, omml, color=NAVY):
        return _add_math_equation(slide, Inches(x_in), Inches(y_in),
                                  Inches(0.98), Inches(0.32), omml,
                                  size_pt=13, color=color)

    DYELLOW = RGBColor(0xB8, 0x86, 0x0B)

    # ---- left panel chrome (static) ---------------------------------
    _add_text(slide, Inches(1.15), Inches(1.53), Inches(4.55),
              Inches(0.35), "Individual Sellers", size=18, bold=True,
              color=NAVY, font="Calibri", align=PP_ALIGN.CENTER)
    _add_arrow(slide, (Inches(1.35), Inches(5.47)),
               (Inches(1.35), Inches(1.87)), color=NAVY, weight_pt=2.0,
               head=True)
    _add_arrow(slide, (Inches(1.35), Inches(5.47)),
               (Inches(5.88), Inches(5.47)), color=NAVY, weight_pt=2.0,
               head=True)
    _add_text(slide, Inches(0.97), Inches(1.55), Inches(0.6),
              Inches(0.3), "P", size=16, bold=True, italic=True,
              color=NAVY, font="Calibri", align=PP_ALIGN.CENTER)
    _add_text(slide, Inches(5.74), Inches(5.51), Inches(0.5),
              Inches(0.3), "q", size=16, bold=True, italic=True,
              color=NAVY, font="Calibri")
    # LMC + LAC (ungrouped, animate as a pair with the header)
    _add_arrow(slide, (Inches(2.48), Inches(5.40)),
               (Inches(5.00), Inches(2.63)), color=FIRM_B_GREEN,
               weight_pt=2.25, head=False)
    _add_text(slide, Inches(4.70), Inches(2.28), Inches(0.85),
              Inches(0.3), "LMC", size=15, bold=True,
              color=FIRM_B_GREEN, font="Calibri")
    _add_bezier_curve(
        slide, (Inches(2.02), Inches(3.25)),
        [((Inches(2.50), Inches(3.90)), (Inches(2.90), Inches(4.23)),
          (Inches(3.60), Inches(4.23))),
         ((Inches(4.50), Inches(4.23)), (Inches(5.10), Inches(3.70)),
          (Inches(5.64), Inches(3.17)))],
        color=GOLD, weight_pt=2.25)
    _add_text(slide, Inches(5.66), Inches(2.91), Inches(0.85),
              Inches(0.3), "LAC", size=14, bold=True, color=GOLD,
              font="Calibri")

    # ---- left-panel groups (Nico's hand groupings, 2026-08-12) ------
    # G: P_Cartel price line + its tick
    a = _add_arrow(slide, (Inches(1.35), Inches(3.31)),
                   (Inches(5.57), Inches(3.31)), color=RED,
                   weight_pt=2.0, head=False)
    b = _tick(0.30, 3.15, _oP('Cartel'), RED)
    _group_shapes(slide, [a, b], "PCartelLine")
    # G: P_Comp dashed line + "d" + its tick
    a = _add_arrow(slide, (Inches(1.35), Inches(4.25)),
                   (Inches(5.57), Inches(4.25)), color=RED,
                   weight_pt=1.5, head=False, dash='dash')
    b = _add_text(slide, Inches(5.59), Inches(4.13), Inches(0.4),
                  Inches(0.3), "d", size=15, bold=True, italic=True,
                  color=RED, font="Calibri")
    c = _tick(0.30, 4.09, _oP('Comp'), RED)
    _group_shapes(slide, [a, b, c], "PCompLine")
    # G: q_Cartel guide + tick
    a = _add_arrow(slide, (Inches(2.74), Inches(3.31)),
                   (Inches(2.74), Inches(5.46)), color=GRAY,
                   weight_pt=1.0, head=False, dash='dash')
    b = _tick(2.32, 5.52, _oq('Cartel'))
    _group_shapes(slide, [a, b], "QCartelGuide")
    # G: q_Dev guide + tick
    a = _add_arrow(slide, (Inches(4.37), Inches(3.31)),
                   (Inches(4.37), Inches(5.46)), color=GRAY,
                   weight_pt=1.0, head=False, dash='dash')
    b = _tick(3.95, 5.52, _oq('Dev'))
    _group_shapes(slide, [a, b], "QDevGuide")
    # G: cartel-profit region + label + arrow
    a = _add_rect(slide, Inches(1.35), Inches(3.31), Inches(1.39),
                  Inches(0.65), PALE_GOLD)
    b = _add_text(slide, Inches(1.04), Inches(2.59), Inches(1.5),
                  Inches(0.55), "Profit\nwith cartel", size=12,
                  bold=True, color=DYELLOW, font="Calibri",
                  align=PP_ALIGN.CENTER)
    c = _add_arrow(slide, (Inches(1.73), Inches(3.02)),
                   (Inches(1.82), Inches(3.60)), color=DYELLOW,
                   weight_pt=1.0, head=True, head_size='sm')
    _group_shapes(slide, [a, b, c], "ProfitRegion")
    # G: deviation extra-profit region + label + arrow
    # green region resized 0.94 -> 0.78 by Nico (2026-08-12)
    a = _add_rect(slide, Inches(2.74), Inches(3.31), Inches(1.63),
                  Inches(0.78), GREEN_FILL)
    b = _add_text(slide, Inches(2.50), Inches(2.66), Inches(1.9),
                  Inches(0.55), "Extra Profit\nwhen deviating",
                  size=12, bold=True, color=FIRM_B_GREEN,
                  font="Calibri", align=PP_ALIGN.CENTER)
    c = _add_arrow(slide, (Inches(3.34), Inches(3.12)),
                   (Inches(3.51), Inches(3.70)), color=FIRM_B_GREEN,
                   weight_pt=1.0, head=True, head_size='sm')
    _group_shapes(slide, [a, b, c], "ExtraProfitRegion")
    # G: q_Comp guide + tick + zero-profit label + arrow
    a = _add_arrow(slide, (Inches(3.52), Inches(4.25)),
                   (Inches(3.52), Inches(5.46)), color=GRAY,
                   weight_pt=1.0, head=False, dash='dash')
    b = _tick(3.10, 5.52, _oq('Comp'))
    c = _add_text(slide, Inches(3.40), Inches(4.72), Inches(1.9),
                  Inches(0.55), "Zero profit\nunder competition",
                  size=12, bold=True, color=RED, font="Calibri",
                  align=PP_ALIGN.CENTER)
    d = _add_arrow(slide, (Inches(3.89), Inches(4.89)),
                   (Inches(3.59), Inches(4.34)), color=RED,
                   weight_pt=1.0, head=True, head_size='sm')
    _group_shapes(slide, [a, b, c, d], "QCompZero")

    # ---- right panel chrome (static) --------------------------------
    _add_text(slide, Inches(7.81), Inches(1.55), Inches(4.55),
              Inches(0.35), "Group/Industry", size=18, bold=True,
              color=NAVY, font="Calibri", align=PP_ALIGN.CENTER)
    _add_arrow(slide, (Inches(7.75), Inches(5.47)),
               (Inches(7.75), Inches(1.87)), color=NAVY, weight_pt=2.0,
               head=True)
    _add_arrow(slide, (Inches(7.75), Inches(5.47)),
               (Inches(12.28), Inches(5.47)), color=NAVY,
               weight_pt=2.0, head=True)
    _add_text(slide, Inches(7.37), Inches(1.55), Inches(0.6),
              Inches(0.3), "P", size=16, bold=True, italic=True,
              color=NAVY, font="Calibri", align=PP_ALIGN.CENTER)
    _add_text(slide, Inches(12.14), Inches(5.51), Inches(0.5),
              Inches(0.3), "Q", size=16, bold=True, italic=True,
              color=NAVY, font="Calibri")
    _add_arrow(slide, (Inches(7.92), Inches(2.42)),
               (Inches(11.63), Inches(5.29)), color=NAVY,
               weight_pt=2.5, head=False)
    _add_text(slide, Inches(11.62), Inches(4.98), Inches(0.4),
              Inches(0.3), "D", size=15, bold=True, color=NAVY,
              font="Calibri")
    _add_arrow(slide, (Inches(8.00), Inches(5.26)),
               (Inches(11.97), Inches(3.63)), color=FIRM_B_GREEN,
               weight_pt=2.5, head=False)
    _add_text(slide, Inches(12.00), Inches(3.32), Inches(0.4),
              Inches(0.3), "S", size=15, bold=True,
              color=FIRM_B_GREEN, font="Calibri")

    # G: MR + label
    a = _add_arrow(slide, (Inches(7.92), Inches(2.42)),
                   (Inches(9.51), Inches(5.31)), color=NAVY,
                   weight_pt=1.75, head=False, dash='dash')
    b = _add_text(slide, Inches(9.53), Inches(5.01), Inches(0.7),
                  Inches(0.3), "MR", size=15, bold=True, color=NAVY,
                  font="Calibri")
    _group_shapes(slide, [a, b], "MRLine")

    def _dot(cx, cy, fill, ln):
        d = Inches(0.11)
        shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     Inches(cx) - d // 2,
                                     Inches(cy) - d // 2, d, d)
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
        if ln is None:
            shp.line.fill.background()
        else:
            shp.line.color.rgb = ln
            shp.line.width = Pt(1.25)
        shp.shadow.inherit = False
        return shp

    # G: cartel outcome (guides + gold dot + ticks)
    a = _add_arrow(slide, (Inches(7.75), Inches(3.39)),
                   (Inches(9.19), Inches(3.39)), color=GRAY,
                   weight_pt=1.0, head=False, dash='dash')
    b = _add_arrow(slide, (Inches(9.19), Inches(3.31)),
                   (Inches(9.19), Inches(5.46)), color=GRAY,
                   weight_pt=1.0, head=False, dash='dash')
    c = _dot(9.185, 3.405, GOLD, NAVY)
    d = _tick(6.72, 3.23, _oP('Cartel'), RED)
    e = _tick(8.74, 5.52, _oQ('Cartel'))
    _group_shapes(slide, [a, b, c, d, e], "CartelOutcome")
    # G: competitive outcome (guides + navy dot + ticks)
    a = _add_arrow(slide, (Inches(7.75), Inches(4.30)),
                   (Inches(10.39), Inches(4.30)), color=GRAY,
                   weight_pt=1.0, head=False, dash='dash')
    b = _add_arrow(slide, (Inches(10.36), Inches(4.34)),
                   (Inches(10.36), Inches(5.46)), color=GRAY,
                   weight_pt=1.0, head=False, dash='dash')
    c = _dot(10.365, 4.285, NAVY, None)
    d = _tick(6.72, 4.14, _oP('Comp'), RED)
    e = _tick(9.91, 5.52, _oQ('C'))
    _group_shapes(slide, [a, b, c, d, e], "CompOutcome")

    # ---- bottom bullets, Nico's extended wording (2026-08-12) -------
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(6.02),
                                  Inches(11.9), Inches(0.82))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    def _mk_run(para, text, *, italic=False, sub=False):
        r = para.add_run()
        r.text = text
        r.font.name = "Calibri"
        r.font.size = Pt(14)
        r.font.italic = italic
        r.font.color.rgb = NAVY
        if sub:
            r._r.get_or_add_rPr().set('baseline', '-25000')
        return r

    lines = [
        [("A monopolist in the market would produce ", 0),
         ("Q", 1), ("cartel", 2), (" ", 1),
         (", and this reduced output leads to a higher market "
          "price ", 0), ("P", 1), ("Cartel", 2)],
        [("For the cartel to work, each firm must produce ", 0),
         ("q", 1), ("Cartel", 2), ("  ", 1),
         ("(that is, each firm must reduce its output)", 0)],
        [("At the cartel price, each firm would like to deviate and "
          "produce more (", 0), ("q", 1), ("Dev", 2),
         (", where ", 0), ("P", 1), ("cartel", 2),
         (" = LMC)", 1)],
    ]
    for i, parts in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if i > 0:
            pPr = para._p.get_or_add_pPr()
            spcBef = ET.SubElement(pPr, qn('a:spcBef'))
            pts = ET.SubElement(spcBef, qn('a:spcPts'))
            pts.set('val', '400')
        for text, kind in parts:
            _mk_run(para, text, italic=(kind >= 1), sub=(kind == 2))
        _set_bullet_char(para, char='\u25aa', color=NAVY, font='Calibri')

    _draw_footer(slide, FOOTER_TEXT, 13)
    _set_notes(slide, (
        "Left panel: the individual cartel member. Under competition "
        "the firm earns zero profit at the LAC minimum. At the cartel "
        "price P_Cartel, producing its quota q_Cartel earns the cartel "
        "profit — but at that price the firm would privately like to "
        "expand to q_Dev, where the price equals its marginal cost "
        "LMC, pocketing the extra profit from deviating. That "
        "temptation is why cartels are unstable.\n"
        "Right panel: the industry. Acting like a monopolist, the "
        "cartel restricts total output to Q_Cartel (where MR crosses "
        "the supply/MC curve) and the price rises above the "
        "competitive level at Q_C."))
    return slide


# --------------------------------------------------------------------------

def slide_14_adm(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_COLLU)
    _draw_action_title(slide, "Collusion: The Case of ADM")

    bullets = [
        ([("Archer Daniels Midland (ADM) involved in organized effort "
           "to raise the price of the animal feed additive ", {}),
          ("lysine", {'link': 'https://en.wikipedia.org/wiki/Lysine'})],
         0),
        ("Unique case with FBI mole", 0),
        ("Taped price-fixing meetings", 1),
        ("Original FBI footage on course webpage (under Module 7 "
         "In-Class Content)", 0),
        ([("Podcast", {'link':
            'https://www.thisamericanlife.org/168/the-fix-is-in'}),
          (", ", {}),
          ("Movie", {'link': 'http://www.imdb.com/title/tt1130080/'}),
          (", Related Podcast: ", {}),
          ("The Poop Cartel", {'link':
            'https://www.npr.org/sections/money/2018/07/25/632444815/'
            'episode-855-the-poop-cartel'})],
         0),
    ]
    _add_hierarchical_bullets(slide, left=MARGIN, top=Inches(1.95),
                              width=Inches(7.6), height=Inches(4.4),
                              items=bullets, size=26, sub_size=24,
                              line_spacing_pts=16)
    # style + wire the hyperlink runs
    link_urls = {
        "lysine": "https://en.wikipedia.org/wiki/Lysine",
        "Podcast": "https://www.thisamericanlife.org/168/the-fix-is-in",
        "Movie": "http://www.imdb.com/title/tt1130080/",
        "The Poop Cartel": "https://www.npr.org/sections/money/2018/07/"
                           "25/632444815/episode-855-the-poop-cartel",
    }
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for p in sh.text_frame.paragraphs:
            for run in p.runs:
                if run.text in link_urls:
                    run.hyperlink.address = link_urls[run.text]
                    run.font.underline = True
                    run.font.color.rgb = RGBColor(0x05, 0x63, 0xC1)

    # hand-edited layout (Nico 2026-08-04): small ADM logo top-right
    # corner, large Informant still
    _add_media_image(slide, "image22.jpeg", left=Inches(11.36),
                     top=Inches(0.49), width=Inches(1.70),
                     height=Inches(1.25), rounded=False, shadow=False)
    _add_media_image(slide, "image23.jpg", left=Inches(8.01),
                     top=Inches(2.12), width=Inches(5.09),
                     height=Inches(2.86), rounded=True, shadow=True)
    _add_text(slide, Inches(8.88), Inches(5.08), Inches(3.4),
              Inches(0.3),
              "Matt Damon in “The Informant!” (2009)",
              size=11, italic=True, color=GRAY, font="Calibri",
              align=PP_ALIGN.CENTER)

    _draw_footer(slide, FOOTER_TEXT, 14)
    _set_notes(slide, (
        "International Price Fixing. It harms consumers, and it is hard "
        "to prove price co-movements are indeed due to collusion.\n"
        "Mark Whitacre, the FBI mole at ADM, recorded the price fixing "
        "meetings. Mention special role of golf on Hawaii.\n"
        "FBI Video: Jump to 2:00. Tyson and ConAgra: major food "
        "producers, customers buying lysine: those that would be most "
        "hurt by price fixing. Actual price fixing at 7:20 – at some "
        "point, they even check the FX rate with the Canadian dollar to "
        "fix the price in a billion-dollar market within a minute.\n"
        "From the Poop Cartel Podcast: “Poop is big business in Dakar, "
        "Senegal. That's because a lot of the homes have septic tanks, "
        "which need to be cleaned out regularly. The best way to clean "
        "your tank is by hiring a proper truck, which arrives with a "
        "vacuum. The problem is that for years, the guys who drove the "
        "poop trucks operated as a cartel, and they squashed competition "
        "and kept prices high. So a lot of people turned to a cheaper "
        "alternative: men who clear out septic tanks by hand, with "
        "shovels and buckets, and basically bury the poop in the street. "
        "It's terrible work and bad for the environment. Today on the "
        "show: how economists got sewage off the streets of Dakar by "
        "making poop-truck drivers compete with each other. In other "
        "words, by breaking up the poop cartel.”"))
    return slide


def slide_16_adm_chart(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_COLLU)
    _draw_action_title(slide, "The ADM Price-Fixing Scandal")

    # lysine price chart (converted from the source EMF, 2026-08-04)
    pic = _add_media_image(slide, "_adm_lysine_chart.png",
                           left=0, top=Inches(1.65),
                           height=Inches(4.85), rounded=False,
                           shadow=True)
    pic.left = int((SLIDE_W - pic.width) // 2)
    _add_text(slide, pic.left, pic.top + pic.height + Inches(0.12),
              pic.width, Inches(0.3),
              "Source: American Economic Association, “Conspiracies "
              "against the public” (aeaweb.org)",
              size=12, italic=True, color=GRAY, font="Calibri",
              align=PP_ALIGN.CENTER)
    _draw_footer(slide, FOOTER_TEXT, 16)
    _set_notes(slide, (
        "The lysine price per pound around the ADM cartel: prices fall "
        "during the price wars, rise when the conspirators meet and "
        "cooperate, and collapse again after the FBI raid. Source: "
        "https://www.aeaweb.org/research/conspiracies-against-the-public"))
    return slide


# --------------------------------------------------------------------------
# Batch B — Bertrand section (29–31, 33–36)
# --------------------------------------------------------------------------

def slide_30_bertrand(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_BERTRAND)
    _draw_action_title(slide, "Bertrand Competition")
    # box narrowed + red "price" emphasis (Nico 2026-08-05)
    _add_hierarchical_bullets(
        slide, left=MARGIN, top=Inches(1.85), width=Inches(9.41),
        height=Inches(3.02),
        items=[
            ([("Bertrand assumed that firms choose a ", {}),
              ("price", {'color': RGBColor(0xFF, 0x00, 0x00)}),
              (" and stand ready to meet all demand at that price",
               {})], 0),
            ("Each firm sets a price that maximizes profits, "
             "conditional on the price set by the other firm", 0),
            ("Contrast this with the Cournot case, where the game was "
             "played on quantities", 0),
        ],
        size=28, sub_size=24, line_spacing_pts=20)
    _draw_footer(slide, FOOTER_TEXT, 30)
    return slide


def slide_31_bertrand_assumptions(prs):
    """Nico's hand design (2026-08-06): full-bleed background photo,
    assumptions card as his scaled group (injected verbatim)."""
    slide = _blank_slide(prs)
    slide.shapes.add_picture(str(SRC_IMG_DIR / "_s30_image43.png"),
                             3658, 0, width=Inches(13.326),
                             height=Inches(7.5))
    _draw_top_bar_tc(slide, TAG_BERTRAND)
    _draw_action_title(slide,
                       "Bertrand Competition (with Identical Goods)")
    _inject_handoff_group(slide, "_handoff_s30_group.xml", id_base=9600)
    _draw_footer(slide, FOOTER_TEXT, 31)
    return slide


def slide_32_concrete_poll(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_BERTRAND)
    _draw_action_title(slide, "Choose Your Concrete Price")

    bullets = [
        ("Two competitors, identical good: Concrete", 0),
        ("Same marginal cost: MC = $120 per cubic yard", 0),
        ("You and your competitor submit a sealed price bid for a huge "
         "project", 0),
        ("What's your price bid?", 0),
    ]
    _add_hierarchical_bullets(slide, left=MARGIN, top=Inches(1.95),
                              width=Inches(7.3), height=Inches(4.2),
                              items=bullets, size=28, sub_size=24,
                              line_spacing_pts=18)
    _add_media_image(slide, "image33.jpeg", left=Inches(8.4),
                     top=Inches(1.9), width=Inches(4.2),
                     rounded=True, shadow=True)       # concrete mixer
    _add_media_image(slide, "image32.jpg", left=Inches(9.55),
                     top=Inches(5.0), width=Inches(1.9),
                     rounded=False, shadow=False)     # $-tug-of-war
    _poll_badge(slide)
    _draw_footer(slide, FOOTER_TEXT, 32)
    _set_notes(slide, (
        "Updated 2026-08-06 (fact-checked): US ready-mix concrete is "
        "quoted per CUBIC YARD, and MC = $120/yd3 matches industry "
        "data - materials (cement, aggregates, admixtures) ran "
        "$89-96/yd3 in 2023-24 plus roughly $15-25/yd3 delivery "
        "(NRMCA benchmarking and monthly surveys). The average US "
        "selling price was about $180/yd3 in 2025, so students should "
        "bid down from around the market price toward the $120 "
        "marginal cost. NOTE: the PollEverywhere answer bands (slide "
        "32) are being updated by Nico to match the new $120 anchor. "
        "Sources: NRMCA Performance Benchmarking Survey (Oct 2024); "
        "concretefinancialinsights.com; US Concrete 10-K."))
    return slide


def slide_34_bertrand_chart(prs):
    """Geometry, labels and tick styling = Nico's hand-edits
    (2026-08-06): tighter chart scale, equation-only demand label,
    'Bertrand Equilibrium' with the down-arrow into the (90,10) point."""
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_BERTRAND)
    _draw_action_title(slide, "Bertrand Competition")

    _add_graphicframe_shadow(slide, Inches(0.42), Inches(1.55),
                             Inches(7.95), Inches(5.3))
    fig = SimpleFig(1.7, 6.25, 6.4, 4.25, 112, 112)
    # shaded regions first
    _add_rect(slide, Inches(1.7), Inches(4.732), Inches(1.714),
              Inches(0.19), RED_FILL)
    _add_rect(slide, Inches(3.414), Inches(4.922), Inches(1.771),
              Inches(0.949), GREEN_FILL)
    _fig_axes(slide, fig)
    _add_text(slide, Inches(1.3), Inches(1.45), Inches(0.7),
              Inches(0.35), "P", size=18, bold=True, italic=True,
              color=NAVY, font="Calibri", align=PP_ALIGN.CENTER)
    _add_text(slide, Inches(8.15), Inches(6.31), Inches(0.6),
              Inches(0.35), "Q", size=18, bold=True, italic=True,
              color=NAVY, font="Calibri")
    # market demand + equation-only label (Nico dropped the bold title)
    _add_arrow(slide, (fig.x(0), fig.y(100)), (fig.x(100), fig.y(0)),
               color=NAVY, weight_pt=2.5, head=False)
    _add_math_equation(slide, Inches(2.187), Inches(2.588),
                       Inches(2.0), Inches(0.4),
                       _omml_run('P') + _omml_text(' = 100 − ')
                       + _oQ('M'), size_pt=15, color=NAVY)
    # MC
    _add_arrow(slide, (fig.x(0), fig.y(10)), (fig.x(105), fig.y(10)),
               color=GRAY, weight_pt=2.5, head=False)
    _add_text(slide, Inches(7.243), Inches(5.511), Inches(0.9),
              Inches(0.3), "MC", size=16, bold=True, color=GRAY,
              font="Calibri", align=PP_ALIGN.CENTER)
    # guides
    _add_arrow(slide, (fig.x(0), fig.y(40)), (fig.x(60), fig.y(40)),
               color=GRAY, weight_pt=1.0, head=False, dash='dash')
    _add_arrow(slide, (fig.x(60), fig.y(40)), (fig.x(60), fig.y(0)),
               color=GRAY, weight_pt=1.0, head=False, dash='dash')
    _add_arrow(slide, (fig.x(0), fig.y(35)), (fig.x(61), fig.y(35)),
               color=GRAY, weight_pt=1.0, head=False, dash='dash')
    _add_arrow(slide, (fig.x(61), fig.y(35)), (fig.x(61), fig.y(0)),
               color=GRAY, weight_pt=1.0, head=False, dash='dash')
    _fig_point(slide, fig, 60, 40, fill=NAVY)
    _fig_point(slide, fig, 61, 35, fill=GOLD, line=NAVY)
    _fig_point(slide, fig, 90, 10, fill=GOLD, line=NAVY, r_in=0.075)
    # region + equilibrium labels
    _add_text(slide, Inches(1.814), Inches(4.712), Inches(1.6),
              Inches(0.26), "Lost by A", size=13, bold=True, color=RED,
              font="Calibri")
    _add_text(slide, Inches(3.586), Inches(5.263), Inches(1.8),
              Inches(0.32), "Gained by A", size=15, bold=True,
              color=FIRM_B_GREEN, font="Calibri")
    # gold action-box callout, same style as slide 25 (2026-08-08)
    _add_rounded_filled_box(
        slide, Inches(6.05), Inches(4.92), Inches(2.2), Inches(0.44),
        "Bertrand Equilibrium", fill=GOLD, text_color=NAVY, size=15,
        bold=True, corner_pct=0.12, shadow=True)
    _add_arrow(slide, (Inches(7.01), Inches(5.3)),
               (Inches(6.88), Inches(5.796)), color=NAVY,
               weight_pt=1.25, head=True)
    # ticks
    _fig_ytick(slide, fig, 100, "$100")
    _fig_ytick(slide, fig, 40, "$40", bold=True)
    _fig_ytick(slide, fig, 35, "$39", bold=True, color=GOLD)
    _fig_ytick(slide, fig, 10, "$10", color=GRAY)
    _fig_xtick(slide, fig, 30, "30")
    # tight tick boxes at Nico's hand geometry (2026-08-12)
    _add_text(slide, Inches(4.88), Inches(6.31), Inches(0.36),
              Inches(0.27), "60", size=16, color=NAVY, font="Calibri",
              align=PP_ALIGN.CENTER)
    _add_text(slide, Inches(5.15), Inches(6.32), Inches(0.41),
              Inches(0.27), "61", size=16, color=GOLD, font="Calibri",
              align=PP_ALIGN.CENTER)
    _fig_xtick(slide, fig, 90, "90", bold=True)
    _fig_xtick(slide, fig, 100, "100")

    # right column
    _add_convention_box(
        slide, Inches(8.6), Inches(1.85), Inches(4.35), Inches(2.5),
        runs=[
            ("Suppose Firms A and B start at a price of $40, where each "
             "supplies Q=30", {'size': 18, 'color': NAVY}),
            ("A believes that B will maintain its price. What will A "
             "do? How will B respond?",
             {'size': 18, 'bold': True, 'color': NAVY, 'newline': True}),
        ],
        fill_rgb=WHITE, border=NAVY, size=18, line_spacing_pct=110)
    _add_rounded_filled_box(
        slide, Inches(8.6), Inches(5.3), Inches(4.35), Inches(1.25),
        "→  A and B will continue to lower their price until they "
        "reach P = MC",
        fill=GOLD, text_color=NAVY, size=19, bold=True, corner_pct=0.12)

    _draw_footer(slide, FOOTER_TEXT, 34)
    _set_notes(slide, (
        "Bertrand assumed that firms choose a price and stand ready to "
        "meet all demand at that price. Each firm sets a price that "
        "maximizes profits holding constant the price set by the other "
        "firm. Contrast this with the Cournot case where the game is "
        "played on quantities.\n"
        "What is the outcome of Bertrand competition? The only "
        "equilibrium is P1 = P2 = MC. If, say P2 >= P1 > MC, P2 will "
        "undercut by a small amount, so there are no such equilibria. "
        "Proof: starting at P1 = P2 = MC, neither firm gains from "
        "raising (sells nothing) or cutting (negative profit). And any "
        "price above MC invites undercutting — the situation in the "
        "figure (both at $40, one drops to $39: it loses the thin strip "
        "on its old sales but gains the whole market).\n"
        "Bertrand with identical goods has a clean solution, but the "
        "conditions are rare. A market that came close: online retail "
        "of computer chips for home-built CPUs — consumers bought the "
        "cheapest chip and prices were approximately marginal cost, "
        "with retailers trying to escape via obfuscation (loss-leader "
        "ads, add-ons). With differentiated goods, equilibrium prices "
        "can stay above marginal cost."))
    return slide

def slide_35_outcomes_comparison(prs):
    """Geometry, annotation and arrow positions = Nico's hand-edits
    (2026-08-06). NO shade on the backing rectangle — deliberate
    exception: the shade would overlap the annotation text that
    extends beyond the figure."""
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_BERTRAND)
    _draw_action_title(slide, "Oligopoly Outcomes: Comparison")

    _add_rect(slide, Inches(0.42), Inches(1.55), Inches(7.6),
              Inches(5.3), WHITE)
    fig = SimpleFig(1.7, 6.25, 6.0, 4.25, 112, 112)
    _fig_axes(slide, fig)
    _add_text(slide, Inches(1.3), Inches(1.45), Inches(0.7),
              Inches(0.35), "P", size=18, bold=True, italic=True,
              color=NAVY, font="Calibri", align=PP_ALIGN.CENTER)
    _add_text(slide, Inches(7.75), Inches(6.31), Inches(0.6),
              Inches(0.35), "Q", size=18, bold=True, italic=True,
              color=NAVY, font="Calibri")
    _add_arrow(slide, (fig.x(0), fig.y(100)), (fig.x(100), fig.y(0)),
               color=NAVY, weight_pt=2.5, head=False)
    _add_math_equation(slide, Inches(1.82), Inches(2.253),
                       Inches(2.0), Inches(0.4),
                       _omml_run('P') + _omml_text(' = 100 − ')
                       + _oQ('M'), size_pt=15, color=NAVY)
    _add_arrow(slide, (fig.x(0), fig.y(100)), (fig.x(50), fig.y(0)),
               color=NAVY, weight_pt=1.75, head=False, dash='dash')
    _add_text(slide, Inches(3.315), Inches(5.219), Inches(0.8),
              Inches(0.32), "MR", size=16, bold=True, color=NAVY,
              font="Calibri")
    _add_arrow(slide, (fig.x(0), fig.y(10)), (fig.x(105), fig.y(10)),
               color=GRAY, weight_pt=2.5, head=False)
    _add_text(slide, Inches(6.896), Inches(5.511), Inches(0.9),
              Inches(0.3), "MC", size=16, bold=True, color=GRAY,
              font="Calibri", align=PP_ALIGN.CENTER)
    for qx, py, fill in [(45, 55, NAVY), (60, 40, NAVY)]:
        _add_arrow(slide, (fig.x(0), fig.y(py)), (fig.x(qx), fig.y(py)),
                   color=GRAY, weight_pt=1.0, head=False, dash='dash')
        _add_arrow(slide, (fig.x(qx), fig.y(py)), (fig.x(qx), fig.y(0)),
                   color=GRAY, weight_pt=1.0, head=False, dash='dash')
        _fig_point(slide, fig, qx, py, fill=fill, line=None)
    # Bertrand horizontal price guide clicks on its OWN (Nico's
    # choreography, 2026-08-13); the rest of the Bertrand cluster is
    # grouped below
    _add_arrow(slide, (fig.x(0), fig.y(10)), (fig.x(90), fig.y(10)),
               color=GRAY, weight_pt=1.0, head=False, dash='dash')
    _fig_ytick(slide, fig, 100, "$100")
    _fig_ytick(slide, fig, 55, "$55", bold=True)
    _fig_ytick(slide, fig, 40, "$40", bold=True)
    _fig_ytick(slide, fig, 10, "$10", color=GRAY)
    for xv in (45, 60):
        _fig_xtick(slide, fig, xv, str(xv))
    # tight '100' tick box (Nico 2026-08-13)
    _add_text(slide, Inches(6.79), Inches(6.31), Inches(0.54),
              Inches(0.27), "100", size=16, color=NAVY, font="Calibri",
              align=PP_ALIGN.CENTER)

    # annotations + arrows at Nico's exact positions/lengths
    _add_text(slide, Inches(4.657), Inches(2.32), Inches(3.093),
              Inches(0.28), "Monopoly Equilibrium", size=16, bold=True,
              color=NAVY, font="Calibri")
    _add_text(slide, Inches(4.657), Inches(2.62), Inches(3.139),
              Inches(0.269), "(also collusive/cartel equilibrium)",
              size=16, color=NAVY, font="Calibri")
    _add_arrow(slide, (Inches(5.015), Inches(3.03)),
               (Inches(4.211), Inches(4.108)), color=NAVY,
               weight_pt=1.0, head=True, head_size='sm')
    _add_text(slide, Inches(5.157), Inches(3.961), Inches(2.275),
              Inches(0.269), "Cournot Equilibrium", size=16, bold=True,
              color=NAVY, font="Calibri")
    _add_arrow(slide, (Inches(5.714), Inches(4.293)),
               (Inches(4.953), Inches(4.695)), color=NAVY,
               weight_pt=1.0, head=True, head_size='sm')
    # Bertrand cluster = ONE hand group (Nico 2026-08-13):
    # vertical guide + gold point + '90' tick + labels + pointer
    _b1 = _add_arrow(slide, (fig.x(90), fig.y(10)), (fig.x(90),
                     fig.y(0)), color=GRAY, weight_pt=1.0, head=False,
                     dash='dash')
    _b2 = _fig_point(slide, fig, 90, 10, fill=GOLD, line=NAVY)
    _b3 = _fig_xtick(slide, fig, 90, "90")
    _b4 = _add_text(slide, Inches(8.02), Inches(4.43), Inches(4.8),
              Inches(0.32), "Bertrand Equilibrium", size=16, bold=True,
              color=NAVY, font="Calibri")
    _b5 = _add_text(slide, Inches(8.02), Inches(4.73), Inches(4.8),
              Inches(0.32), "(same as Perfect Competition", size=16,
              color=NAVY, font="Calibri")
    _b6 = _add_text(slide, Inches(8.02), Inches(5.03), Inches(4.8),
              Inches(0.32), "when both firms have same MC)", size=16,
              color=NAVY, font="Calibri")
    _b7 = _add_arrow(slide, (Inches(7.915), Inches(4.65)),
               (Inches(6.546), Inches(5.811)), color=NAVY,
               weight_pt=1.0, head=True, head_size='sm')
    _group_shapes(slide, [_b1, _b2, _b3, _b4, _b5, _b6, _b7],
                  "BertrandCluster")

    _draw_footer(slide, FOOTER_TEXT, 35)
    _set_notes(slide, (
        "You can use this graph to understand the welfare consequences "
        "of these different forms of competition.\n"
        "Producer welfare ranks: Perfect Competition = Bertrand < "
        "Cournot < Monopoly. Consumer welfare ranks: Perfect "
        "Competition = Bertrand > Cournot > Monopoly. Total surplus "
        "ranks: Perfect Competition = Bertrand > Cournot > Monopoly.\n"
        "These last two comparisons (consumer welfare and total surplus "
        "decrease with market power) are the reason why the Federal "
        "Trade Commission cares about market power."))
    return slide

def slide_36_outcomes_math(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_BERTRAND)
    _draw_action_title(slide, "Comparison of Outcomes: Math (for your reference)")

    # bold + nudged to y 1.58 by Nico (2026-08-13)
    _eqbox = _add_math_equation(
        slide, Inches(4.1), Inches(1.58), Inches(5.1), Inches(0.5),
        _omml_run('P') + _omml_text(' = 100 − ') + _omml_run('Q')
        + _omml_text('   and   MC = 10'),
        size_pt=20, color=GRAY)
    # OMML ignores a:rPr@b — bold comes from m:sty (bi = bold
    # italic for variables, b = bold upright for text/numbers)
    _MQ = '{%s}%%s' % M_NS
    for _mr in _eqbox._element.iter(_MQ % 'r'):
        _mrPr = _mr.find(_MQ % 'rPr')
        if _mrPr is None:
            _mrPr = ET.Element(_MQ % 'rPr')
            _mr.insert(0, _mrPr)
        _sty = _mrPr.find(_MQ % 'sty')
        if _sty is None:
            _sty = ET.SubElement(_mrPr, _MQ % 'sty')
            _sty.set(_MQ % 'val', 'bi')
        else:
            _sty.set(_MQ % 'val',
                     'b' if _sty.get(_MQ % 'val') == 'p' else 'bi')
        for _ar in _mr.findall('{%s}rPr' % A_NS):
            _ar.set('b', '1')

    segs = [
        ("text", "Monopoly:", {'bold': True, 'size': 24}),
        ("break", None, {}),
        ("omml", _omml_text('TR = ') + _omml_run('P') + _omml_text('·')
         + _omml_run('Q') + _omml_text(' = (100 − ') + _omml_run('Q')
         + _omml_text(')·') + _omml_run('Q') + _omml_text(' = 100')
         + _omml_run('Q') + _omml_text(' − ')
         + _omml_sup(_omml_run('Q'), _omml_text('2')), {'size': 20}),
        ("break", None, {}),
        ("omml", _omml_text('MR = 100 − 2') + _omml_run('Q'),
         {'size': 20}),
        ("break", None, {}),
        ("omml", _omml_text('MR = MC   →   100 − 2') + _omml_run('Q')
         + _omml_text(' = 10'), {'size': 20}),
        ("break", None, {}),
        ("text", "➜  Q = 45,  P = 55", {'bold': True, 'size': 20,
                                        'color': RED}),
        ("break", None, {}), ("break", None, {}), ("break", None, {}),
        ("text", "Cournot:  ", {'bold': True, 'size': 24}),
        ("text", "Q = 60,  P = 40", {'bold': True, 'size': 20,
                                     'color': RED}),
        ("break", None, {}), ("break", None, {}), ("break", None, {}),
        ("text", "Perfect competition / Bertrand:", {'bold': True,
                                                     'size': 24}),
        ("break", None, {}),
        ("omml", _omml_run('P') + _omml_text(' = MC   →   100 − ')
         + _omml_run('Q') + _omml_text(' = 10'), {'size': 20}),
        ("break", None, {}),
        ("text", "➜  Q = 90,  P = 10", {'bold': True, 'size': 20,
                                        'color': RED}),
    ]
    _add_mixed_textbox(slide, Inches(0.63), Inches(2.27), Inches(9.5),
                       Inches(4.65), segs, default_size=20,
                       default_color=NAVY)
    # gold separator rules with shade above each case (Nico 2026-08-06)
    for sy in (2.2, 4.52, 5.53):   # sep3 up 0.09, Nico 2026-08-13
        rule = _add_rect(slide, Inches(0.63), Inches(sy), Inches(9.5),
                         Inches(0.028), GOLD)
        _add_drop_shadow(rule)
    _draw_footer(slide, FOOTER_TEXT, 36)
    _add_video_link_box(slide,
                        "Practice Video “Cournot Competition Math”",
                        size=16)
    _set_notes(slide, (
        "The algebra behind the comparison chart, all with the same "
        "market demand P = 100 − Q and MC = 10. The monopolist (or "
        "cartel) sets MR = MC and produces 45 at a price of 55. The "
        "Cournot duopoly ends up at 60 and a price of 40 — between "
        "monopoly and competition. Bertrand price competition drives "
        "the price all the way down to marginal cost: 90 units at a "
        "price of 10, the perfectly competitive outcome."))
    return slide


def slide_37_overview(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_BERTRAND)
    _draw_action_title(slide, "Overview: Cournot vs. Bertrand")

    col_w = Inches(5.95)
    gap = Inches(0.45)
    left1 = (SLIDE_W - 2 * col_w - gap) // 2
    left2 = left1 + col_w + gap
    _add_rounded_filled_box(slide, left1, Inches(1.85), col_w,
                            Inches(0.6), "Cournot: quantity competition",
                            fill=NAVY, text_color=WHITE, size=22,
                            bold=True, corner_pct=0.15)
    _add_rounded_filled_box(slide, left2, Inches(1.85), col_w,
                            Inches(0.6), "Bertrand: price competition",
                            fill=NAVY, text_color=WHITE, size=22,
                            bold=True, corner_pct=0.15)
    _add_bulleted_list(slide, left=left1 + Inches(0.1),
                       top=Inches(2.7), width=col_w - Inches(0.2),
                       height=Inches(0.9),
                       items=["Firms set capacity strategically"],
                       size=22, line_spacing_pts=12)
    _add_bulleted_list(slide, left=left2 + Inches(0.1),
                       top=Inches(2.7), width=col_w - Inches(0.2),
                       height=Inches(0.9),
                       items=["Firms set price strategically"],
                       size=22, line_spacing_pts=12)
    _add_hierarchical_bullets(
        slide, left=MARGIN + Inches(0.2), top=Inches(3.9),
        width=Inches(12.2), height=Inches(2.6),
        items=[
            ("There isn't a “correct” and an “incorrect” model, they "
             "just make different assumptions", 0),
            ("Often Bertrand (price) competition in the short run", 1),
            ("Cournot (quantity) competition in the long run (e.g., by "
             "setting up production capacities)", 1),
        ],
        size=24, sub_size=22, line_spacing_pts=14)
    _draw_footer(slide, FOOTER_TEXT, 37)
    _set_notes(slide, (
        "Bertrand with identical goods has a nice solution but the "
        "conditions that give rise to this outcome are rare in the real "
        "world. A market that was close was the online market for "
        "computer chips for consumers that like to build their own "
        "CPUs. Consumers bought the cheapest chip. Researchers found "
        "that prices were approx. marginal cost. An interesting wrinkle "
        "is that computer chip retailers attempted to differentiate "
        "their goods through obfuscation: either by advertising an "
        "inferior chip as a loss leader and then redirecting the "
        "consumer to their website or by offering product add-ons. In a "
        "Bertrand model with differentiated goods, there is the "
        "potential for the equilibrium price to be above marginal "
        "cost."))
    return slide


# --------------------------------------------------------------------------
# Batch C — Oligopoly with differentiated goods (38–46). Firm colors from
# the source deck: Airbus / Firm A = red, Boeing / Firm B = green.
# --------------------------------------------------------------------------

def _oPA():
    return _omml_sub(_omml_run('P'), _omml_text('A'))


def _oPB():
    return _omml_sub(_omml_run('P'), _omml_text('B'))


def slide_39_diff_products(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_DIFF)
    _draw_action_title(slide,
                       "Bertrand Competition with Differentiated Products")
    # underline emphasis + box height per Nico's hand edits (2026-08-07)
    _add_hierarchical_bullets(
        slide, left=MARGIN, top=Inches(1.85), width=RULE_W,
        height=Inches(3.59),   # auto-fit after Nico's bullet deletion
        items=[
            ([("So far: We assumed that all firms produce ", {}),
              ("identical products", {'underline': True})], 0),
            ([("More common situation: ", {}),
              ("Differentiated", {'underline': True}),
              (" products", {})], 0),
            ("Pepsi and Coke", 1),
            ("Airbus A320 vs. Boeing 737 (medium-range)", 1),
            ("Airbus A350 vs. Boeing 787 (long-range)", 1),
            ("Setting now:", 0),
            ([("2 firms producing ", {}),
              ("differentiated", {'underline': True}),
              (" products", {})], 1),
        ],
        size=28, sub_size=24, line_spacing_pts=14)
    _draw_footer(slide, FOOTER_TEXT, 39)
    return slide


def slide_40_diff_assumptions(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_DIFF)
    _draw_action_title(slide,
                       "Bertrand Competition (with Differentiated Goods)")
    # red "Differentiated" in the TITLE (Nico's hand edit 2026-08-07)
    for sh in slide.shapes:
        if (sh.has_text_frame and
                sh.text_frame.text.startswith("Bertrand Competition (")):
            para = sh.text_frame.paragraphs[0]
            for r in list(para.runs):
                r._r.getparent().remove(r._r)
            for t, col in [("Bertrand Competition (with ", NAVY),
                           ("Differentiated", RED),
                           (" Goods)", NAVY)]:
                r = para.add_run()
                r.text = t
                r.font.name = "Calibri"
                r.font.size = Pt(32)
                r.font.bold = True
                r.font.color.rgb = col
            break
    box_w = Inches(10.6)
    left = (SLIDE_W - box_w) // 2
    box = _add_convention_box(
        slide, left, Inches(2.7), box_w, Inches(2.6),
        runs=[
            ("Assumptions:", {'size': 26, 'bold': True, 'color': NAVY}),
            ("▪  Firms do not sell identical products. They sell ",
             {'size': 24, 'color': NAVY, 'newline': True}),
            ("differentiated", {'size': 24, 'color': RED}),
            (" products", {'size': 24, 'color': NAVY}),
            ("▪  Each firm faces its own demand function!",
             {'size': 24, 'color': NAVY, 'newline': True}),
            ("▪  Each firm ", {'size': 24, 'color': NAVY,
                              'newline': True}),
            ("simultaneously", {'size': 24, 'color': NAVY}),
            (" chooses the price at which it sells its product",
             {'size': 24, 'color': NAVY}),
        ],
        size=24, line_spacing_pct=135)
    _add_drop_shadow(box)   # match slide 21 (Nico 2026-08-04)
    # underline "simultaneously" (hand edit)
    for sh in slide.shapes:
        if sh.has_text_frame:
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    if run.text == "simultaneously":
                        run.font.underline = True
    _draw_footer(slide, FOOTER_TEXT, 40)
    _set_notes(slide, (
        "The three assumptions for Bertrand with differentiated "
        "goods: differentiated products, an own demand function for "
        "each firm, and prices chosen simultaneously. Why do we need "
        "simultaneity? It means neither firm can observe the rival's "
        "price before setting its own, and neither can commit to "
        "moving first — so the equilibrium is a pair of mutual best "
        "responses rather than a leader-follower outcome. It also "
        "answers the natural question about the price war: the second "
        "mover cannot simply wait and undercut at the last moment, "
        "because in this setting there is no last mover. In Part 2 we "
        "deliberately break exactly this assumption — commitment is "
        "what changes when one player CAN move first."))
    return slide


def slide_41_diff_setup(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_DIFF)
    _draw_action_title(slide, "Oligopoly with Differentiated Goods")

    _add_hierarchical_bullets(
        slide, left=MARGIN, top=Inches(1.7), width=Inches(8.3),
        height=Inches(1.5),
        items=[
            ("Oligopoly on long-range aircraft:", 0),
            ([("Airbus A350 ", {}),
              ("(Firm A)", {'bold': True, 'color': FIRM_A_RED})], 1),
            ([("Boeing 787 ", {}),
              ("(Firm B)", {'bold': True, 'color': FIRM_B_GREEN})], 1),
        ],
        size=26, sub_size=24, line_spacing_pts=8)
    # picture sizes/positions hand-tweaked by Nico (2026-08-08)
    _add_media_image(slide, "image35.png", left=Inches(8.67),
                     top=Inches(1.465), width=Inches(4.08),
                     rounded=True, shadow=True)     # A350
    _add_media_image(slide, "image36.png", left=Inches(7.708),
                     top=Inches(3.35), width=Inches(5.042),
                     rounded=True, shadow=True)     # 787

    _add_text(slide, MARGIN, Inches(3.55), Inches(5.0), Inches(0.45),
              "Interdependent demand:", size=26, color=NAVY,
              font="Calibri")
    eq_a = (_omml_sub(_omml_run('Q', color=FIRM_A_RED),
                      _omml_text('A', color=FIRM_A_RED))
            + _omml_text(' = 400 − 2')
            + _omml_sub(_omml_run('P', color=FIRM_A_RED),
                        _omml_text('A', color=FIRM_A_RED))
            + _omml_text(' + ')
            + _omml_sub(_omml_run('P', color=ACC3_50),
                        _omml_text('B', color=ACC3_50)))
    eq_b = (_omml_sub(_omml_run('Q', color=ACC3_50),
                      _omml_text('B', color=ACC3_50))
            + _omml_text(' = 400 − 2')
            + _omml_sub(_omml_run('P', color=ACC3_50),
                        _omml_text('B', color=ACC3_50))
            + _omml_text(' + ')
            + _omml_sub(_omml_run('P', color=FIRM_A_RED),
                        _omml_text('A', color=FIRM_A_RED)))
    _add_math_equation(slide, Inches(1.7), Inches(4.05), Inches(4.4),
                       Inches(0.5), eq_a, size_pt=24, color=NAVY)
    _add_math_equation(slide, Inches(1.7), Inches(4.6), Inches(4.4),
                       Inches(0.5), eq_b, size_pt=24, color=NAVY)
    # substitutes illustration per the ORIGINAL slide (Nico
    # 2026-08-07): gold circle around the "+ P_A" term of Q_B's demand,
    # one arrow, "substitutes" label
    # circle/line/label geometry = Nico's hand edits (2026-08-08)
    oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.783),
                                  Inches(4.56), Inches(1.047),
                                  Inches(0.62))
    oval.fill.background()
    oval.line.color.rgb = GOLD
    oval.line.width = Pt(2.25)
    oval.shadow.inherit = False
    _add_text(slide, Inches(6.475), Inches(4.47), Inches(1.7),
              Inches(0.35), "substitutes", size=16, italic=True,
              bold=True, color=GOLD, font="Calibri")
    _add_arrow(slide, (Inches(6.41), Inches(4.66)),
               (Inches(5.79), Inches(4.785)), color=GOLD,
               weight_pt=3.0, head=False)
    nbox = _add_convention_box(
        slide, Inches(6.92), Inches(5.2), Inches(4.27), Inches(0.75),
        prefix="Note:  ",
        body="Now we have 2 demand functions! (Differentiated products)",
        size=18)
    _add_drop_shadow(nbox)

    _add_text(slide, MARGIN, Inches(5.55), Inches(6.5), Inches(0.45),
              "Same marginal cost (all in $Million):", size=26,
              color=NAVY, font="Calibri")
    _add_math_equation(slide, Inches(1.7), Inches(6.05), Inches(4.2),
                       Inches(0.5),
                       _omml_sub(_omml_text('MC', color=FIRM_A_RED),
                                 _omml_text('A', color=FIRM_A_RED))
                       + _omml_text(' = ')
                       + _omml_sub(_omml_text('MC', color=ACC3_50),
                                   _omml_text('B', color=ACC3_50))
                       + _omml_text(' = 200'),
                       size_pt=24, color=NAVY)
    _draw_footer(slide, FOOTER_TEXT, 41)
    _add_video_link_box(slide,
                        "Optional Advanced Practice Video: “Bertrand "
                        "Competition with Differentiated Goods – Math”",
                        size=13)
    return slide


def _airbus_reaction_slide(prs, page_num, *, roman, pb, dem_int,
                           opt_q, opt_p, prev=False):
    """Airbus residual-demand chart, arranged like the ORIGINAL slides
    (Nico 2026-08-07): NO shade (would overlap the text), assumption +
    response as red-bordered callout pair top-right, circle on the
    cross-price term of the demand label, MRA at the MR foot."""
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_DIFF)
    _draw_action_title(
        slide, f"Reaction Function: Airbus' Reaction to Boeing's "
               f"Price ({roman})")

    fig = SimpleFig(2.4, 6.35, 6.3, 4.5, 880, 460)
    _fig_axes(slide, fig)
    # P-label nudged by hand per slide (Nico 2026-08-14)
    _plx, _ply = (1.93, 1.52) if prev else (1.89, 1.52)
    _add_text(slide, Inches(_plx), Inches(_ply), Inches(0.7),
              Inches(0.35), "P", size=18, bold=True, italic=True,
              color=NAVY, font="Calibri", align=PP_ALIGN.CENTER)
    _add_text(slide, Inches(8.85), Inches(6.41), Inches(0.6),
              Inches(0.35), "Q", size=18, bold=True, italic=True,
              color=NAVY, font="Calibri")

    color_new = FIRM_A_RED if not prev else FIRM_B_GREEN
    text_new = FIRM_A_RED if not prev else ACC3_50
    if prev:
        # previous demand at PB=200 (red) + pointer label + shift arrow
        _add_arrow(slide, (fig.x(0), fig.y(300)), (fig.x(600),
                   fig.y(0)), color=FIRM_A_RED, weight_pt=2.0,
                   head=False)
        _add_mixed_textbox(
            slide, Inches(6.981), Inches(5.914), Inches(2.917),
            Inches(0.32),
            [("text", "Previous demand at ", {'size': 14, 'bold': True,
                                              'color': FIRM_A_RED}),
             ("omml", _omml_sub(_omml_run('P', color=FIRM_A_RED),
                                _omml_text('B', color=FIRM_A_RED))
              + _omml_text(' = 200', color=FIRM_A_RED), {'size': 14})],
            default_size=14, default_color=FIRM_A_RED)
        _add_arrow(slide, (Inches(6.961), Inches(6.113)),
                   (Inches(6.581), Inches(6.263)), color=FIRM_A_RED,
                   weight_pt=2.0, head=False)
        # green arrow marking the outward demand shift
        _add_arrow(slide, (Inches(6.287), Inches(6.059)),
                   (Inches(6.772), Inches(5.486)), color=FIRM_B_GREEN,
                   weight_pt=2.75, head=True)
    # demand + MR at the assumed PB
    _add_arrow(slide, (fig.x(0), fig.y(dem_int)),
               (fig.x(dem_int * 2), fig.y(0)),
               color=color_new, weight_pt=2.5, head=False)
    _mr_line = _add_arrow(slide, (fig.x(0), fig.y(dem_int)),
               (fig.x(dem_int), fig.y(0)),
               color=color_new, weight_pt=2.0, head=False, dash='dash')
    mra_pos = (4.746, 5.58) if prev else (4.215, 5.8)
    _mra_box = _add_math_equation(slide, Inches(mra_pos[0]),
                       Inches(mra_pos[1]),
                       Inches(0.95), Inches(0.36), _oMR('A'),
                       size_pt=15, color=text_new)
    # demand label + circled cross-price term + side note
    if prev:
        lbl_x, lbl_y = 2.995, 2.182
        seg2 = [("omml", _omml_text('    = 400 − 0.5·', color=text_new)
                 + _omml_sub(_omml_run('Q', color=FIRM_A_RED),
                             _omml_text('A', color=FIRM_A_RED)),
                 {'size': 14})]
        oval_geo = (5.667, 2.387, 0.336, 0.42)
        tag_txt, tag_pos = "400", (6.205, 1.954)
        tagline = ((6.155, 2.204), (5.955, 2.414))
        noteline = ((7.719, 5.114), (7.066, 5.561))
        note = [("text", "Higher demand for Airbus because ",
                 {'size': 17, 'bold': True, 'color': ACC3_50}),
                ("omml", _omml_sub(_omml_run('P', color=ACC3_50),
                                   _omml_text('B', color=ACC3_50)),
                 {'size': 17}),
                ("text", " has increased", {'size': 17, 'bold': True,
                                            'color': ACC3_50})]
        note_pos = (7.795, 4.81)
    else:
        lbl_x, lbl_y = 2.814, 2.862
        seg2 = [("omml", _omml_text('    = 200 − 0.5·', color=text_new)
                 + _omml_sub(_omml_run('Q', color=FIRM_A_RED),
                             _omml_text('A', color=FIRM_A_RED))
                 + _omml_text(' + ', color=text_new)
                 + _omml_text('100', color=ACC3_50), {'size': 14}),
                ("break", None, {}),
                ("omml", _omml_text('    = 300 − 0.5·', color=text_new)
                 + _omml_sub(_omml_run('Q', color=FIRM_A_RED),
                             _omml_text('A', color=FIRM_A_RED)),
                 {'size': 14})]
        oval_geo = (5.5, 3.057, 0.26, 0.4)
        tag_txt, tag_pos = "200", (5.835, 2.62)
        tagline = ((5.835, 2.851), (5.685, 3.061))
        noteline = ((6.066, 3.515), (5.75, 3.358))
        note = [("text", "A higher price ", {'size': 17, 'bold': True,
                                             'color': ACC3_50}),
                ("omml", _omml_sub(_omml_run('P', color=ACC3_50),
                                   _omml_text('B', color=ACC3_50)),
                 {'size': 17}),
                ("text", " increases demand for Airbus",
                 {'size': 17, 'bold': True, 'color': ACC3_50})]
        note_pos = (6.065, 3.36)
    segs = [("text", "Airbus' (inverse) Demand",
             {'bold': True, 'size': 15, 'color': text_new}),
            ("break", None, {}),
            ("omml", _omml_sub(_omml_run('P', color=FIRM_A_RED),
                               _omml_text('A', color=FIRM_A_RED))
             + _omml_text(' = 200 − 0.5·', color=text_new)
             + _omml_sub(_omml_run('Q', color=FIRM_A_RED),
                         _omml_text('A', color=FIRM_A_RED))
             + _omml_text(' + 0.5·', color=ACC3_50)
             + _omml_sub(_omml_run('P', color=ACC3_50),
                         _omml_text('B', color=ACC3_50)),
             {'size': 14}),
            ("break", None, {})] + seg2
    _add_mixed_textbox(slide, Inches(lbl_x), Inches(lbl_y),
                       Inches(3.1), Inches(1.15), segs,
                       default_size=14, default_color=text_new)
    oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(oval_geo[0]),
                                  Inches(oval_geo[1]),
                                  Inches(oval_geo[2]),
                                  Inches(oval_geo[3]))
    oval.fill.background()
    oval.line.color.rgb = ACC3_50
    oval.line.width = Pt(2.0)
    oval.shadow.inherit = False
    _tag_box = _add_text(slide, Inches(tag_pos[0]), Inches(tag_pos[1]),
              Inches(0.9), Inches(0.35), tag_txt, size=16, italic=True,
              color=ACC3_50, font="Calibri")
    if prev:
        # circle + tag = one hand group on slide 43 (Nico 2026-08-14)
        _group_shapes(slide, [oval, _tag_box], "CircleTag")
    _add_arrow(slide, (Inches(tagline[0][0]),
                       Inches(tagline[0][1])),
               (Inches(tagline[1][0]), Inches(tagline[1][1])),
               color=ACC3_50, weight_pt=2.0, head=False)
    _add_mixed_textbox(slide, Inches(note_pos[0]),
                       Inches(note_pos[1]), Inches(3.4), Inches(0.8),
                       note, default_size=17, default_color=ACC3_50)
    # short connector from the note toward the circle (Nico's added
    # line, 2026-08-08)
    _add_arrow(slide, (Inches(noteline[0][0]), Inches(noteline[0][1])),
               (Inches(noteline[1][0]), Inches(noteline[1][1])),
               color=ACC3_50, weight_pt=2.0, head=False)
    # MC line
    _add_arrow(slide, (fig.x(0), fig.y(200)), (fig.x(850), fig.y(200)),
               color=GRAY, weight_pt=2.5, head=False)
    mc_pos = (8.056, 4.033) if prev else (7.95, 4.469)
    _add_text(slide, Inches(mc_pos[0]), Inches(mc_pos[1]),
              Inches(0.8), Inches(0.269), "MC", size=16, bold=True,
              italic=True, color=GRAY, font="Calibri",
              align=PP_ALIGN.CENTER)
    # optimum
    _add_arrow(slide, (fig.x(opt_q), fig.y(0)), (fig.x(opt_q),
               fig.y(opt_p)), color=GRAY, weight_pt=1.0, head=False,
               dash='sysDot')
    _add_arrow(slide, (fig.x(0), fig.y(opt_p)), (fig.x(opt_q),
               fig.y(opt_p)), color=GRAY, weight_pt=1.0, head=False,
               dash='sysDot')
    _fig_point(slide, fig, opt_q, 200, fill=NAVY)
    _fig_point(slide, fig, opt_q, opt_p, fill=GOLD, line=NAVY)
    _fig_ytick(slide, fig, dem_int, str(dem_int))
    _fig_ytick(slide, fig, opt_p, str(opt_p), bold=True)
    _fig_ytick(slide, fig, 200, "200", color=GRAY)
    _fig_xtick(slide, fig, opt_q, str(opt_q), bold=True)
    _mr_tick = _fig_xtick(slide, fig, dem_int, str(dem_int))
    _fig_xtick(slide, fig, dem_int * 2, str(dem_int * 2))
    if not prev:
        # MR line + label + its foot tick = one hand group on slide 42
        _group_shapes(slide, [_mr_line, _mra_box, _mr_tick], "MRGroup")

    # paired red-bordered callouts (original layout, top-right)
    box1 = _add_outlined_box(slide, Inches(9.65), Inches(1.72),
                             Inches(3.3), Inches(0.9), "",
                             line=FIRM_A_RED, fill=WHITE, line_w=2.0,
                             rounded=False, shadow=False)
    p1 = box1.text_frame.paragraphs[0]
    for t, opts in [("Airbus assumes that Boeing charges ",
                     dict(bold=False, color=FIRM_A_RED)),
                    ("P", dict(bold=True, italic=True, color=ACC3_50)),
                    ("B", dict(bold=True, italic=True, color=ACC3_50,
                               sub=True)),
                    (f" = {pb}", dict(bold=True, italic=True,
                                      color=ACC3_50))]:
        r = p1.add_run()
        r.text = t
        r.font.name = "Calibri"
        r.font.size = Pt(18)
        r.font.bold = opts.get('bold', False)
        r.font.italic = opts.get('italic', False)
        r.font.color.rgb = opts['color']
        if opts.get('sub'):
            r._r.get_or_add_rPr().set('baseline', '-25000')
    box2 = _add_outlined_box(slide, Inches(9.65), Inches(2.82),
                             Inches(3.3), Inches(0.95), "",
                             line=FIRM_A_RED, fill=WHITE, line_w=2.0,
                             rounded=False, shadow=False)
    p2 = box2.text_frame.paragraphs[0]
    for t, sub in [("→ Airbus' optimal price response is ", False),
                   ("P", False), ("A", True), (f" = {opt_p}", False)]:
        r = p2.add_run()
        r.text = t
        r.font.name = "Calibri"
        r.font.size = Pt(18)
        r.font.bold = True
        r.font.italic = t in ("P", "A")
        r.font.color.rgb = FIRM_A_RED
        if sub:
            r._r.get_or_add_rPr().set('baseline', '-25000')

    _draw_footer(slide, FOOTER_TEXT, page_num)
    _set_notes(slide, (
        f"Same two-step logic as in the Cournot case, but now in "
        f"prices. Step 1: fix Boeing's price at PB={pb} and substitute "
        f"it into Airbus' demand, giving the inverse demand "
        f"PA = {dem_int} − 0.5·QA. Step 2: maximize profit as usual "
        f"(MR = MC with MC = 200): the optimal quantity is {opt_q} and "
        f"the optimal price response is PA = {opt_p}. Note the "
        f"difference to Cournot: price responds POSITIVELY to price "
        f"increases by the other firm."))
    return slide

def slide_42_airbus_i(prs):
    return _airbus_reaction_slide(prs, 42, roman="I", pb=200,
                                  dem_int=300, opt_q=100, opt_p=250)


def slide_43_airbus_ii(prs):
    return _airbus_reaction_slide(prs, 43, roman="II", pb=400,
                                  dem_int=400, opt_q=200, opt_p=300,
                                  prev=True)


def slide_44_equilibrium_prices(prs):
    """Arranged like the ORIGINAL slide + Nico's hand groups
    (2026-08-14): seven clusters (A/B reactions, starting point,
    reaction points I/II, the 267 guides, the equilibrium callout),
    each ONE group; choreography reveals the points BEFORE the lines."""
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_DIFF)
    _draw_action_title(slide, "Equilibrium Prices")

    fig = SimpleFig(2.75, 6.5, 6.4, 4.75, 470, 350)
    _fig_axes(slide, fig)
    ttl = _add_mixed_textbox(
        slide, Inches(0.538), Inches(3.486), Inches(3.0), Inches(0.32),
        [("text", "Airbus' Price (", {'bold': True, 'italic': True,
                                      'size': 16, 'color': FIRM_A_RED}),
         ("omml", _omml_sub(_omml_run('P', color=FIRM_A_RED),
                            _omml_text('A', color=FIRM_A_RED)),
          {'size': 16, 'color': FIRM_A_RED}),
         ("text", ")", {'bold': True, 'italic': True, 'size': 16,
                        'color': FIRM_A_RED})],
        default_size=16, default_color=FIRM_A_RED)
    ttl.rotation = 270
    _add_mixed_textbox(
        slide, Inches(5.0), Inches(6.85), Inches(3.0), Inches(0.35),
        [("text", "Boeing's Price (", {'bold': True, 'italic': True,
                                       'size': 16, 'color': ACC3_50}),
         ("omml", _omml_sub(_omml_run('P', color=ACC3_50),
                            _omml_text('B', color=ACC3_50)),
          {'size': 16, 'color': ACC3_50}),
         ("text", ")", {'bold': True, 'italic': True, 'size': 16,
                        'color': ACC3_50})],
        default_size=16, default_color=ACC3_50)

    # G0: Airbus reaction line + label + its intercept tick (200 on P_A)
    _a1 = _add_arrow(slide, (fig.x(0), fig.y(200)), (fig.x(450),
                     fig.y(312.5)), color=FIRM_A_RED, weight_pt=2.5,
                     head=False)
    _a2 = _add_mixed_textbox(
        slide, Inches(8.96), Inches(2.058), Inches(2.9), Inches(0.8),
        [("text", "Airbus' reaction function", {'bold': True,
                                                'size': 14,
                                                'color': FIRM_A_RED}),
         ("break", None, {}),
         ("omml", _omml_sub(_omml_run('P', color=FIRM_A_RED),
                            _omml_text('A', color=FIRM_A_RED))
          + _omml_text(' = 200 + 0.25·', color=FIRM_A_RED)
          + _omml_sub(_omml_run('P', color=ACC3_50),
                      _omml_text('B', color=ACC3_50)), {'size': 14})],
        default_size=14, default_color=FIRM_A_RED)
    _a3 = _fig_ytick(slide, fig, 200, "200", color=GRAY)
    _group_shapes(slide, [_a1, _a2, _a3], "ReactionA")

    # G1: Boeing reaction line + label
    _b1 = _add_arrow(slide, (fig.x(200), fig.y(0)), (fig.x(287.5),
                     fig.y(350)), color=FIRM_B_GREEN, weight_pt=2.5,
                     head=False)
    _b2 = _add_mixed_textbox(
        slide, Inches(6.27), Inches(1.41), Inches(2.9), Inches(0.8),
        [("text", "Boeing's reaction function", {'bold': True,
                                                 'size': 14,
                                                 'color': ACC3_50}),
         ("break", None, {}),
         ("omml", _omml_sub(_omml_run('P', color=ACC3_50),
                            _omml_text('B', color=ACC3_50))
          + _omml_text(' = 200 + 0.25·', color=ACC3_50)
          + _omml_sub(_omml_run('P', color=FIRM_A_RED),
                      _omml_text('A', color=FIRM_A_RED)), {'size': 14})],
        default_size=14, default_color=ACC3_50)
    _group_shapes(slide, [_b1, _b2], "ReactionB")

    # G2: starting point (tick + dotted guide + dotted card)
    _s1 = _fig_ytick(slide, fig, 150, "150", color=GRAY)
    _s2 = _add_arrow(slide, (fig.x(0), fig.y(150)), (fig.x(237.5),
                     fig.y(150)), color=GRAY, weight_pt=1.5, head=True,
                     dash='sysDot')
    box = _add_outlined_box(slide, Inches(8.718), Inches(5.0),
                            Inches(3.35), Inches(0.9), "",
                            line=GRAY, fill=WHITE, line_w=2.0,
                            rounded=False, shadow=False)
    _dash_shape_line(box, 'sysDot')
    pbx = box.text_frame.paragraphs[0]
    for t, sub in [("Starting point: Airbus contemplates charging ",
                    False), ("P", False), ("A", True),
                   (" = 150", False)]:
        r = pbx.add_run()
        r.text = t
        r.font.name = "Calibri"
        r.font.size = Pt(15)
        r.font.italic = True
        r.font.color.rgb = GRAY
        if sub:
            r._r.get_or_add_rPr().set('baseline', '-25000')
    _group_shapes(slide, [_s1, _s2, box], "Starting")

    # G3 / G4: reaction points I and II (guides + dot + ticks + label)
    _rp_groups = []
    for (bx, ay, lbl, lpos, apos) in [
            (200, 250, "Reaction point (I)", (6.62, 4.62),
             ((6.59, 4.65), (5.55, 3.22))),
            (400, 300, "Reaction point (II)", (8.718, 3.323),
             ((8.618, 3.395), (8.24, 2.489)))]:
        g = []
        g.append(_add_arrow(slide, (fig.x(0), fig.y(ay)), (fig.x(bx),
                 fig.y(ay)), color=ACC6_75, weight_pt=1.5, head=False,
                 dash='sysDot'))
        g.append(_add_arrow(slide, (fig.x(bx), fig.y(ay)), (fig.x(bx),
                 fig.y(0)), color=ACC6_75, weight_pt=1.5, head=False,
                 dash='sysDot'))
        g.append(_fig_point(slide, fig, bx, ay, fill=GOLD, line=NAVY,
                            r_in=0.065))
        g.append(_fig_ytick(slide, fig, ay, str(ay), color=ACC6_75))
        g.append(_fig_xtick(slide, fig, bx, str(bx), color=ACC6_75))
        g.append(_add_text(slide, Inches(lpos[0]), Inches(lpos[1]),
                 Inches(2.2), Inches(0.3), lbl, size=14, bold=True,
                 italic=True, color=ACC6_75, font="Calibri"))
        g.append(_add_arrow(slide, (Inches(apos[0][0]),
                 Inches(apos[0][1])), (Inches(apos[1][0]),
                 Inches(apos[1][1])), color=ACC6_75, weight_pt=2.0,
                 head=False))
        _rp_groups.append(_group_shapes(slide, g, "ReactionPoint"))

    # convergence staircase (3 separate clicks, ungrouped)
    for sp_, ep_ in [((5.958, 4.464), (5.958, 3.0)),
                     ((5.96, 3.042), (6.31, 3.042)),
                     ((6.3, 3.047), (6.3, 2.891))]:
        _add_arrow(slide, (Inches(sp_[0]), Inches(sp_[1])),
                   (Inches(ep_[0]), Inches(ep_[1])), color=GRAY,
                   weight_pt=1.5, head=True, dash='sysDot')

    # G5: the 267-267 dashed guides + both ticks
    _c1 = _add_arrow(slide, (fig.x(0), fig.y(267)), (fig.x(267),
                     fig.y(267)), color=NAVY, weight_pt=1.75,
                     head=False, dash='dash')
    _c2 = _add_arrow(slide, (fig.x(267), fig.y(267)), (fig.x(267),
                     fig.y(0)), color=NAVY, weight_pt=1.75,
                     head=False, dash='dash')
    _c3 = _fig_ytick(slide, fig, 267, "267", bold=True)
    _c4 = _fig_xtick(slide, fig, 267, "267", bold=True)
    _group_shapes(slide, [_c1, _c2, _c3, _c4], "TwoSixtySeven")

    # G6: equilibrium point + gold callout + pointer arrow
    _e1 = _fig_point(slide, fig, 267, 267, fill=GOLD, line=NAVY,
                     r_in=0.075)
    _e2 = _add_arrow(slide, (Inches(7.0), Inches(3.68)),
                     (Inches(6.5), Inches(3.0)), color=NAVY,
                     weight_pt=1.75, head=True)
    _e3 = _add_rounded_filled_box(
        slide, Inches(7.05), Inches(3.62), Inches(2.05), Inches(0.73),
        "", fill=GOLD, text_color=NAVY, size=14, bold=True,
        corner_pct=0.12, shadow=True)
    _e4 = _add_mixed_textbox(
        slide, Inches(7.1), Inches(3.66), Inches(1.95), Inches(0.65),
        [("text", "Equilibrium", {'bold': True, 'size': 15}),
         ("break", None, {}),
         ("omml", _omml_sub(_omml_run('P', color=FIRM_A_RED),
                            _omml_text('A', color=FIRM_A_RED))
          + _omml_text(' = ')
          + _omml_sub(_omml_run('P', color=ACC3_50),
                      _omml_text('B', color=ACC3_50))
          + _omml_text(' = 267'), {'size': 14})],
        default_size=14, default_color=NAVY)
    _group_shapes(slide, [_e1, _e2, _e3, _e4], "Equilibrium")

    _draw_footer(slide, FOOTER_TEXT, 44)
    _set_notes(slide, (
        "Price responds POSITIVELY to price increases by the other "
        "firm. Note that this is the opposite of Cournot, where Firm A "
        "wants to produce less if firm B produces more. The two "
        "reaction functions — Airbus' PA = 200 + 0.25·PB and Boeing's "
        "mirror image — intersect at the equilibrium PA = PB = 267. "
        "Starting from PA = 150, the contemplated responses move both "
        "firms toward the equilibrium."))
    return slide


def slide_45_bruinlearn(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_DIFF)
    _draw_action_title(slide, "Play This Game on Bruinlearn")
    _add_hierarchical_bullets(
        slide, left=MARGIN, top=Inches(2.1), width=Inches(12.3),
        height=Inches(3.2),
        items=[
            ("Under Module 7 / In-Class Content", 0),
            ("You can choose Boeing's price and observe Airbus's "
             "reaction (price)", 0),
            ("You also see the profits of both companies", 1),
            ("You can also change the demand functions (make payoffs "
             "asymmetric)", 0),
        ],
        size=28, sub_size=24, line_spacing_pts=18)
    _add_outlined_box(
        slide, Inches(3.67), Inches(5.14), Inches(6.0), Inches(0.85),
        "BruinLearn  →  Module 7 / In-Class Content",
        line=GOLD, text_color=NAVY, size=20, bold=True, line_w=1.75,
        rounded=True, shadow=True, corner_pct=0.20)
    _draw_footer(slide, FOOTER_TEXT, 45)
    return slide


def slide_46_takeaways(prs):
    s = make_content_bulleted(
        prs, page_num=46, section_tag=TAG_DIFF,
        title="Take-Aways: Oligopoly with Differentiated Goods",
        bullets=[
            ("The two firms produce differentiated products, so each "
             "has its own demand curve", 0),
            ("However, each firm's demand curve is affected by the "
             "other firm's price", 0),
            ([("When the competitor raises its price, the optimal "
               "response is for the other firm to ", {}),
              ("also raise", {'underline': True}),
              (" its price", {})], 0),
            ("In equilibrium, neither firm has an incentive to "
             "deviate", 0),
            ("If MC are different, the optimal prices will be "
             "different for the two firms", 0),
            ([("Note: Math on this problem is ", {}),
              ("optional", {'underline': True}),
              (" (will not be in final exam)", {})], 0),
        ],
        size=24, sub_size=22, line_spacing_pts=12)
    _add_video_link_box(s,
                        "Optional Advanced Practice Video: “Bertrand "
                        "Competition with Differentiated Goods – Math”",
                        size=13)
    return s


def slide_47_concluding_discussion(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_DIFF)
    _draw_action_title(slide, "Oligopoly: Concluding Discussion")
    _add_hierarchical_bullets(
        slide, left=MARGIN, top=Inches(2.2), width=Inches(12.3),
        height=Inches(3.2),
        items=[
            ("If you could pick one type of Oligopoly competition for "
             "your firm, which would you choose? Why?", 0),
            ("In which oligopoly case can you improve your product?", 0),
            ("What is the effect of improving your product?", 1),
        ],
        size=28, sub_size=24, line_spacing_pts=22)
    _add_discussion_break(slide)
    _draw_footer(slide, FOOTER_TEXT, 47)
    return slide


# --------------------------------------------------------------------------
# Batch D — Game theory: concepts (48–53)
# --------------------------------------------------------------------------

def slide_49_what_is_game_theory(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_GT)
    _draw_action_title(slide, "What Is Game Theory?")
    _add_hierarchical_bullets(
        slide, left=MARGIN, top=Inches(1.9), width=Inches(8.3),
        height=Inches(4.6),
        items=[
            ("Branch of economics concerned with interdependent "
             "situations", 0),
            ("The outcome of a participant's choice depends critically "
             "on the actions of other participants", 1),
            ("Game theory has transformed not only economics, but also "
             "the study of war, of business strategy, and of "
             "biological systems", 0),
            ("Father of Game Theory: John Nash", 0),
            ("1994 Nobel Prize in Economics", 1),
            ("Inspired the movie “A Beautiful Mind”", 1),
        ],
        size=24, sub_size=22, line_spacing_pts=14)
    # geometry + z-order = Nico's hand edits (2026-08-08): medal
    # first so it sits BEHIND the Nash portrait
    _add_media_image(slide, "image38.jpeg", left=Inches(10.933),
                     top=Inches(1.95), height=Inches(1.6),
                     rounded=False, shadow=False)  # Nobel medal
    pic = _add_media_image(slide, "image37.jpeg", left=Inches(8.813),
                           top=Inches(1.401), height=Inches(2.949),
                           rounded=True, shadow=True)   # John Nash
    _add_text(slide, Inches(9.038), Inches(4.43),
              Inches(1.95), Inches(0.28), "John Nash (1928–2015)",
              size=11, italic=True, color=GRAY, font="Calibri",
              align=PP_ALIGN.CENTER)
    _add_media_image(slide, "image39.jpeg", left=Inches(6.466),
                     top=Inches(4.62), height=Inches(2.775),
                     rounded=False, shadow=True)   # Beautiful Mind DVD
    _draw_footer(slide, FOOTER_TEXT, 49)
    _set_notes(slide, (
        "In biology, game theory is used to explain how organisms "
        "behave when their success depends on the actions of others. "
        "It helps clarify why an organism might cooperate, compete, "
        "form alliances, or adopt particular signaling or mating "
        "strategies. Classic examples include the evolution of "
        "cooperation, predator–prey interactions, and stable ratios of "
        "aggressive vs. passive behavioral types within a species."))
    return slide


def slide_50_warmup_game(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_GT)
    _draw_action_title(slide, "A Warm-Up Game")
    _add_hierarchical_bullets(
        slide, left=MARGIN, top=Inches(1.9), width=Inches(8.0),
        height=Inches(4.6),
        items=[
            ("I brought 5 yellow and 5 red cards", 0),
            ("But this is not about soccer…", 0),
            ("Instead:", 0),
            ("I distribute the 5 yellow cards to students", 1),
            ("Suppose I get $100 from the Dean for each red-yellow "
             "pair that I take back to her after class", 1),
            ("You can ‘sell’ your yellow card back to me", 1),
            ("What's your asking price for your yellow card?", 1),
        ],
        size=24, sub_size=22, line_spacing_pts=12)
    _add_media_image(slide, "image40.jpeg", left=Inches(9.0),
                     top=Inches(2.3), width=Inches(3.7),
                     rounded=True, shadow=True)
    _draw_footer(slide, FOOTER_TEXT, 50)
    return slide


def slide_51_beauty_contest(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_GT)
    _draw_action_title(slide, "Strategic Thinking: The Beauty Contest")
    _add_hierarchical_bullets(
        slide, left=MARGIN, top=Inches(1.9), width=Inches(6.854),
        height=Inches(3.256),
        items=[
            ("Prominent game in newspapers in the 1930s:", 0),
            ([("Get a reward if you select the face that is the most "
               "popular among others (not the one that ", {}),
              ("you", {'underline': True}),
              (" find most attractive)", {})], 1),
            ("John Maynard Keynes (1936) used it as analogy", 0),
        ],
        size=26, sub_size=24, line_spacing_pts=18)
    # picture size/position = Nico's hand edits (2026-08-08)
    pic = _add_media_image(slide, "image41.jpeg", left=Inches(7.4),
                           top=Inches(3.51), width=Inches(5.58),
                           height=Inches(2.988), rounded=True,
                           shadow=True)
    _add_text(slide, Inches(7.4), Inches(6.577),
              Inches(5.58), Inches(0.372), "John Maynard Keynes",
              size=11, italic=True, color=GRAY, font="Calibri",
              align=PP_ALIGN.CENTER)
    _draw_footer(slide, FOOTER_TEXT, 51)
    return slide


def slide_52_numerical_beauty(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_GT)
    _draw_action_title(slide,
                       "Strategic Thinking: Numerical Beauty Contest")
    _add_hierarchical_bullets(
        slide, left=MARGIN, top=Inches(2.4), width=Inches(12.3),
        height=Inches(2.8),
        items=[
            ("Numerical application:", 0),
            ("All participants are asked to pick a number between 0 "
             "and 100", 1),
            ("Winner: the number that is closest to 2/3 of the average "
             "of all numbers submitted", 1),
        ],
        size=28, sub_size=24, line_spacing_pts=18)
    _poll_badge(slide)
    _draw_footer(slide, FOOTER_TEXT, 52)
    return slide


def slide_54_key_concepts(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_GT)
    _draw_action_title(slide, "Key Concepts of Game Theory")
    CONCEPT_BLUE = RGBColor(0x00, 0x70, 0xC0)
    _add_hierarchical_bullets(
        slide, left=MARGIN, top=Inches(1.95), width=Inches(12.4),
        height=Inches(4.6),
        items=[
            ("Strategic behavior by all participants", 0),
            ([("Optimization: maximize ", {}),
              ("individual", {'italic': True}),
              (" payoff", {})], 1),
            ("To study strategic behavior:", 0),
            ([("Best response", {'bold': True, 'color': CONCEPT_BLUE}),
              (": optimal decision given what the other player does",
               {})], 1),
            ([("Dominant strategy", {'bold': True,
                                     'color': CONCEPT_BLUE}),
              (": optimal action no matter what the other player does",
               {})], 1),
            ([("Nash equilibrium", {'bold': True,
                                    'color': CONCEPT_BLUE}),
              (": «intersection of best responses»", {})], 1),
            ("Position from which neither player would find it "
             "advantageous to deviate (e.g., Cournot equilibrium)", 2,
             {'size': 22}),
        ],
        size=26, sub_size=24, line_spacing_pts=14)
    _draw_footer(slide, FOOTER_TEXT, 54)
    return slide


# --------------------------------------------------------------------------
# Batch E — classic games (55–61, 63–68, 70–73, 75, 79–84)
# --------------------------------------------------------------------------

def _pd_matrix(slide, *, dim=None, nash=False, arrows_p1=False,
               arrows_p2=False, top=Inches(2.5)):
    """The prisoner's-dilemma matrix used on slides 55-58."""
    anchors = _add_payoff_matrix(
        slide, left=Inches(4.55), top=top, cell_w=Inches(2.3),
        cell_h=Inches(1.05),
        row_player="Prisoner 1", col_player="Prisoner 2",
        row_strats=["Silent", "Betray"], col_strats=["Silent", "Betray"],
        payoffs=[[("-1", "-1"), ("-10", "0")],
                 [("0", "-10"), ("-8", "-8")]],
        caption=True, nash_cells=[(1, 1)] if nash else (), dim=dim)
    if arrows_p1:
        for c in (0, 1):
            _br_arrow(slide, anchors[(0, c, 'row')],
                      anchors[(1, c, 'row')], ROW_BLUE)
            _br_circle(slide, anchors[(1, c, 'row')], ROW_BLUE)
    if arrows_p2:
        for r in (0, 1):
            _br_arrow(slide, anchors[(r, 0, 'col')],
                      anchors[(r, 1, 'col')], GOLD)
            _br_circle(slide, anchors[(r, 1, 'col')], GOLD)
    return anchors


def slide_56_pd_setup(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_GAMES)
    _draw_action_title(slide, "Prisoner's Dilemma: Setup and Payoffs")
    _pd_matrix(slide, top=Inches(2.55))
    _add_convention_box(
        slide, Inches(0.5), Inches(5.6), Inches(4.9), Inches(1.15),
        runs=[("Interpretation of payoffs:", {'size': 18,
                                              'underline': True, 'color': NAVY}),
              ("Negative values, because they indicate years in prison",
               {'size': 17, 'color': NAVY, 'newline': True})],
        size=17)
    _draw_footer(slide, FOOTER_TEXT, 56)
    _set_notes(slide, (
        "We start with the payoff table. Payoffs to the player on the "
        "left (Prisoner 1) are blue, and shown on the left. Payoffs to "
        "the player on top of the matrix (Prisoner 2) are gold, and "
        "shown on the right."))
    return slide


def slide_57_pd_p1(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_GAMES)
    _draw_action_title(slide,
                       "Prisoner's Dilemma: Strategy of Prisoner 1")
    _pd_matrix(slide, dim='row', arrows_p1=True, top=Inches(2.55))
    _add_convention_box(
        slide, Inches(0.5), Inches(5.6), Inches(4.9), Inches(1.3),
        runs=[("Interpretation of arrows:", {'size': 18,
                                             'underline': True, 'color': NAVY}),
              ("Prisoner 1 picks his best response to Prisoner 2's "
               "(hypothetical) choices",
               {'size': 17, 'color': NAVY, 'newline': True})],
        size=17)
    _draw_footer(slide, FOOTER_TEXT, 57)
    _set_notes(slide, (
        "This slide analyzes the strategy of Prisoner 1. Here we blind "
        "out the payoffs for Prisoner 2 because Prisoner 1 only cares "
        "about his own payoffs when deciding on a strategy. The arrows "
        "for Prisoner 1 (in general, for the player to the left of the "
        "matrix) are drawn vertically."))
    return slide


def slide_58_pd_p2(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_GAMES)
    _draw_action_title(slide,
                       "Prisoner's Dilemma: Strategy of Prisoner 2")
    _pd_matrix(slide, dim='col', arrows_p2=True, top=Inches(2.55))
    _add_convention_box(
        slide, Inches(0.5), Inches(5.6), Inches(4.9), Inches(1.3),
        runs=[("Interpretation of arrows:", {'size': 18,
                                             'underline': True, 'color': NAVY}),
              ("Prisoner 2 picks his best response to Prisoner 1's "
               "(hypothetical) choices",
               {'size': 17, 'color': NAVY, 'newline': True})],
        size=17)
    _draw_footer(slide, FOOTER_TEXT, 58)
    _set_notes(slide, (
        "And now we blind out the payoffs for Prisoner 1 because "
        "Prisoner 2 only cares about his own payoffs when deciding on "
        "a strategy. The arrows for Prisoner 2 (in general, the player "
        "on top of the matrix) are drawn horizontally."))
    return slide


def slide_59_pd_equilibrium(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_GAMES)
    _draw_action_title(slide, "Prisoner's Dilemma: Equilibrium")
    anchors = _pd_matrix(slide, nash=True, arrows_p1=True,
                         arrows_p2=True, top=Inches(2.55))
    _add_outlined_box(
        slide, Inches(10.35), Inches(3.3), Inches(2.65), Inches(0.62),
        "Nash Equilibrium", line=GOLD, text_color=NAVY, size=18,
        bold=True, line_w=1.75, rounded=True, shadow=True,
        corner_pct=0.25)
    cx, cy = anchors[(1, 1, 'cell')]
    _add_arrow(slide, (Inches(10.75), Inches(3.92)),
               (cx + Inches(1.0), cy - Inches(0.2)), color=GOLD,
               weight_pt=1.75, head=True)
    _add_convention_box(
        slide, Inches(0.4), Inches(5.6), Inches(5.1), Inches(1.35),
        runs=[("Both players' ", {'size': 17, 'color': NAVY}),
              ("dominant strategy", {'size': 17, 'italic': True,
                                     'color': NAVY}),
              (" is to betray, so individual optimization produces a "
               "“socially” sub-optimal outcome",
               {'size': 17, 'color': NAVY})],
        size=17)
    _draw_footer(slide, FOOTER_TEXT, 59)
    _set_notes(slide, (
        "In the Nash Equilibrium, neither player has the incentive to "
        "deviate: If Prisoner 1 betrays, then Prisoner 2 also wants to "
        "betray, and vice-versa."))
    return slide


def slide_60_solve_dilemma(prs):
    s = make_content_bulleted(
        prs, page_num=60, section_tag=TAG_GAMES,
        title="Ways to Solve the Dilemma",
        bullets=[
            ("Communication:", 0),
            ("What would a lawyer do?", 1),
            ("External solutions:", 0),
            ("What would the Mafia do?", 1),
            ("Commitment", 1),
            ("Moral rules (guilt/“conscience”)", 1),
            ("Repetition", 0),
            ("Podcast (next slide)", 1),
        ],
        size=28, sub_size=24, line_spacing_pts=14)
    return s


def slide_61_podcast(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_GAMES)
    _draw_action_title(slide, "Podcast on Prisoner's Dilemma")
    _add_hierarchical_bullets(
        slide, left=MARGIN, top=Inches(2.0), width=Inches(12.3),
        height=Inches(1.6),
        items=[
            ("What's the experiment that Bob Axelrod ran in the "
             "1980s?", 0),
            ("What was the winning strategy?", 0),
        ],
        size=28, sub_size=24, line_spacing_pts=16)
    _add_media_image(slide, "image45.png", left=Inches(3.4),
                     top=Inches(3.9), width=Inches(6.5),
                     rounded=False, shadow=True)
    _draw_footer(slide, FOOTER_TEXT, 61)
    return slide


def slide_62_fair_play(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_GAMES)
    _draw_action_title(slide, "Fair Play vs. “Dirty” Play")
    anchors = _add_payoff_matrix(
        slide, left=Inches(4.85), top=Inches(2.6), cell_w=Inches(2.3),
        cell_h=Inches(1.05),
        row_player="Team 1", col_player="Team 2",
        row_strats=["Play fair", "Play dirty"],
        col_strats=["Play fair", "Play dirty"],
        payoffs=[[("0.5", "0.5"), ("0.3", "0.7")],
                 [("0.7", "0.3"), ("0.5", "0.5")]],
        caption=True, nash_cells=[(1, 1)])
    for c in (0, 1):
        _br_arrow(slide, anchors[(0, c, 'row')], anchors[(1, c, 'row')],
                  ROW_BLUE)
        _br_circle(slide, anchors[(1, c, 'row')], ROW_BLUE)
    for r in (0, 1):
        _br_arrow(slide, anchors[(r, 0, 'col')], anchors[(r, 1, 'col')],
                  GOLD)
        _br_circle(slide, anchors[(r, 1, 'col')], GOLD)
    _add_outlined_box(
        slide, Inches(10.55), Inches(3.4), Inches(2.55), Inches(0.62),
        "Nash Equilibrium", line=GOLD, text_color=NAVY, size=17,
        bold=True, line_w=1.75, rounded=True, shadow=True,
        corner_pct=0.25)
    cx, cy = anchors[(1, 1, 'cell')]
    _add_arrow(slide, (Inches(11.0), Inches(4.02)),
               (cx + Inches(1.0), cy - Inches(0.2)), color=GOLD,
               weight_pt=1.75, head=True)
    _add_convention_box(
        slide, Inches(0.35), Inches(5.75), Inches(3.4), Inches(1.1),
        runs=[("Interpretation of payoffs:", {'size': 16,
                                              'underline': True, 'color': NAVY}),
              ("Teams' odds of winning", {'size': 16, 'color': NAVY,
                                          'newline': True})],
        size=16)
    _add_rounded_filled_box(
        slide, Inches(4.15), Inches(6.15), Inches(6.5), Inches(0.7),
        "Each team's dominant strategy is to play dirty",
        fill=GOLD, text_color=NAVY, size=20, bold=True,
        corner_pct=0.15)
    _draw_footer(slide, FOOTER_TEXT, 62)
    return slide


def slide_64_collective_action(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_GAMES)
    _draw_action_title(slide, "General Lesson: Collective Action Problem")
    _add_hierarchical_bullets(
        slide, left=MARGIN, top=Inches(1.7), width=Inches(12.4),
        height=Inches(1.6),
        items=[
            ("Acting selfishly can lead to socially suboptimal "
             "outcomes", 0),
            ("Also holds in situations with more than 2 players", 1),
            ("Known as “tragedy of the commons”", 1),
        ],
        size=26, sub_size=24, line_spacing_pts=8)
    # ORIGINAL deck's commons image (Nico 2026-08-09), with the
    # deck-standard rounded corners + shade
    pic = _add_media_image(slide, "_s63_commons.png", left=0,
                           top=Inches(3.35), height=Inches(2.95),
                           rounded=True, shadow=True)
    pic.left = int((SLIDE_W - pic.width) // 2)
    _add_text(slide, MARGIN, Inches(6.5), Inches(12.4), Inches(0.5),
              "▪  Stark contrast with “invisible hand” principle "
              "(Adam Smith)", size=26, color=NAVY, font="Calibri")
    _draw_footer(slide, FOOTER_TEXT, 64)
    return slide


def slide_65_chicken(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_GAMES)
    _draw_action_title(slide, "The Game of Chicken")
    _add_text(slide, MARGIN, Inches(1.32), Inches(9.5), Inches(0.35),
              "You and your arch-enemy drive towards each other on a "
              "narrow street", size=15, italic=True, color=GRAY,
              font="Calibri")
    anchors = _add_payoff_matrix(
        slide, left=Inches(4.55), top=Inches(2.6), cell_w=Inches(2.3),
        cell_h=Inches(1.0),
        row_player="Arch enemy", col_player="You",
        row_strats=["avoid", "drive on"],
        col_strats=["avoid", "drive on"],
        payoffs=[[("0", "0"), ("-10", "10")],
                 [("10", "-10"), ("-100", "-100")]],
        caption=True, nash_cells=[(0, 1), (1, 0)])
    # best responses: avoid iff the other drives on
    _br_arrow(slide, anchors[(0, 0, 'row')], anchors[(1, 0, 'row')],
              ROW_BLUE)
    _br_circle(slide, anchors[(1, 0, 'row')], ROW_BLUE)
    _br_arrow(slide, anchors[(1, 1, 'row')], anchors[(0, 1, 'row')],
              ROW_BLUE)
    _br_circle(slide, anchors[(0, 1, 'row')], ROW_BLUE)
    _br_arrow(slide, anchors[(0, 0, 'col')], anchors[(0, 1, 'col')],
              GOLD)
    _br_circle(slide, anchors[(0, 1, 'col')], GOLD)
    _br_arrow(slide, anchors[(1, 1, 'col')], anchors[(1, 0, 'col')],
              GOLD)
    _br_circle(slide, anchors[(1, 0, 'col')], GOLD)
    _add_rounded_filled_box(
        slide, Inches(10.55), Inches(3.35), Inches(2.45), Inches(0.7),
        "2 Nash Equilibria", fill=GOLD, text_color=NAVY, size=17,
        bold=True, corner_pct=0.25)
    link = _add_text(slide, Inches(9.55), Inches(1.75), Inches(3.3),
                     Inches(0.35), "Similar game in “Rebel Without a "
                     "Cause”", size=14, italic=True, color=NAVY,
                     font="Calibri", align=PP_ALIGN.RIGHT)
    for p in link.text_frame.paragraphs:
        for run in p.runs:
            run.hyperlink.address = ("https://www.youtube.com/"
                                     "watch?v=BGtEp7zFdrc")
    _add_hierarchical_bullets(
        slide, left=MARGIN, top=Inches(5.85), width=Inches(12.3),
        height=Inches(1.2),
        items=[
            ("Real world examples:", 0),
            ("Entry into a market that can only sustain one firm "
             "profitably (natural monopoly)", 1),
            ("The Cold War (next slide)", 1),
        ],
        size=20, sub_size=18, line_spacing_pts=4)
    _draw_footer(slide, FOOTER_TEXT, 65)
    return slide


def slide_66_mad(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_GAMES)
    _draw_action_title(slide, "The Cold War and M.A.D.")
    _add_hierarchical_bullets(
        slide, left=MARGIN, top=Inches(1.7), width=Inches(9.2),
        height=Inches(1.2),
        items=[
            ("Why did none of the two super-powers use nukes?", 0),
            ("Mutually Assured Destruction (M.A.D.)", 1),
        ],
        size=24, sub_size=22, line_spacing_pts=6)
    anchors = _add_payoff_matrix(
        slide, left=Inches(4.75), top=Inches(3.6), cell_w=Inches(2.45),
        cell_h=Inches(1.0),
        row_player="USA", col_player="USSR",
        row_strats=["No\nnukes", "Use\nnukes"],
        col_strats=["Don't use\nnuclear weapons", "Use nuclear\nweapons"],
        payoffs=[[("0", "0"), ("-20", "10")],
                 [("10", "-20"), ("-100", "-100")]],
        caption=True, nash_cells=[(0, 1), (1, 0)])
    _br_arrow(slide, anchors[(0, 0, 'row')], anchors[(1, 0, 'row')],
              ROW_BLUE)
    _br_circle(slide, anchors[(1, 0, 'row')], ROW_BLUE)
    _br_arrow(slide, anchors[(1, 1, 'row')], anchors[(0, 1, 'row')],
              ROW_BLUE)
    _br_circle(slide, anchors[(0, 1, 'row')], ROW_BLUE)
    _br_arrow(slide, anchors[(0, 0, 'col')], anchors[(0, 1, 'col')],
              GOLD)
    _br_circle(slide, anchors[(0, 1, 'col')], GOLD)
    _br_arrow(slide, anchors[(1, 1, 'col')], anchors[(1, 0, 'col')],
              GOLD)
    _br_circle(slide, anchors[(1, 0, 'col')], GOLD)
    _add_rounded_filled_box(
        slide, Inches(10.6), Inches(4.35), Inches(2.45), Inches(0.7),
        "2 Nash Equilibria", fill=GOLD, text_color=NAVY, size=17,
        bold=True, corner_pct=0.25)
    _add_media_image(slide, "image48.png", left=Inches(0.55),
                     top=Inches(3.2), height=Inches(3.4),
                     rounded=True, shadow=True)   # Titan II launch (tall)
    _add_media_image(slide, "image47.png", left=Inches(9.7),
                     top=Inches(1.5), width=Inches(3.2),
                     rounded=True, shadow=True)   # Soviet warhead (wide)
    _draw_footer(slide, FOOTER_TEXT, 66)
    _set_notes(slide, (
        "Picture left: US Titan II nuclear missile. Picture right: "
        "Soviet Union nuclear warhead."))
    return slide


def slide_67_group_work(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_GAMES)
    _draw_action_title(slide, "Group Work")
    _add_hierarchical_bullets(
        slide, left=MARGIN, top=Inches(1.7), width=Inches(12.3),
        height=Inches(1.2),
        items=[
            ("How does the payoff table change under M.A.D.?", 0),
            ([("“Doomsday device” (also explained by ", {}),
              ("Dr. Strangelove", {'underline': True,
                                   'color': RGBColor(0x05, 0x63, 0xC1)}),
              (")", {})], 1),
        ],
        size=24, sub_size=22, line_spacing_pts=6)
    # full answer UNDER the cover box (original technique, Nico
    # 2026-08-09): payoffs, impossible-cell notes, best-response
    # arrows and the outcome oval are all hidden by the gradient
    # cover, which the instructor deletes/moves live in class
    anchors = _add_payoff_matrix(
        slide, left=Inches(4.75), top=Inches(3.6), cell_w=Inches(2.45),
        cell_h=Inches(1.0),
        row_player="USA", col_player="USSR",
        row_strats=["No\nnukes", "Use\nnukes"],
        col_strats=["Don't use\nnuclear weapons", "Use nuclear\nweapons"],
        payoffs=[[("0", "0"), (None, None)],
                 [(None, None), ("-100", "-100")]],
        caption=True,
        cell_texts={(0, 1): "(Cell impossible due to M.A.D.)",
                    (1, 0): "(Cell impossible due to M.A.D.)"})
    # small outcome oval around just the "0, 0" pair (original style;
    # a full-cell oval would poke out from under the cover box)
    _br_circle(slide, anchors[(0, 0, 'cell')], GOLD, w_in=1.4,
               h_in=0.55, weight_pt=2.75)
    for c in (0, 1):
        _br_arrow(slide, anchors[(1, c, 'row')],
                  anchors[(0, c, 'row')], ROW_BLUE)
    for r in (0, 1):
        _br_arrow(slide, anchors[(r, 1, 'col')],
                  anchors[(r, 0, 'col')], GOLD)
    _add_cover_box(slide, Inches(4.78), Inches(3.63), Inches(4.84),
                   Inches(1.94))
    mixed = [("text", "The world came very close to destruction:  ",
              {'size': 20}),
             ("text", "Documentary", {'size': 20, 'bold': True,
                                      'underline': True,
                                      'color': RGBColor(0x05, 0x63,
                                                        0xC1)})]
    _add_mixed_textbox(slide, MARGIN, Inches(6.35), Inches(8.0),
                       Inches(0.5), mixed, default_size=20,
                       default_color=NAVY)
    # hyperlinks
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for p in sh.text_frame.paragraphs:
            for run in p.runs:
                if run.text == "Dr. Strangelove":
                    run.hyperlink.address = ("https://www.imdb.com/"
                                             "title/tt0057012/")
                elif run.text == "Documentary":
                    run.hyperlink.address = ("https://www.youtube.com/"
                                             "watch?v=4VPY2SgyG5w")
    _add_discussion_break(slide)
    _draw_footer(slide, FOOTER_TEXT, 67)
    _set_notes(slide, (
        "See here for an illustration using a decision tree: "
        "https://www.spaceship.com.au/learn/game-theory-cold-war/\n"
        "Documentary: “The Man Who Saved the World”: Vasili Arkhipov."))
    return slide


def slide_68_commitment(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_GAMES)
    _draw_action_title(slide, "Commitment and Strategy")
    _add_hierarchical_bullets(
        slide, left=MARGIN, top=Inches(1.65), width=Inches(12.3),
        height=Inches(0.9),
        items=[
            ([("Note that M.A.D. represents a ", {}),
              ("commitment", {'underline': True}),
              (" (to retaliate)", {})], 0),
        ],
        size=24, sub_size=22, line_spacing_pts=6)
    anchors = _add_payoff_matrix(
        slide, left=Inches(4.75), top=Inches(3.35), cell_w=Inches(2.45),
        cell_h=Inches(1.0),
        row_player="USA", col_player="USSR",
        row_strats=["No\nnukes", "Use\nnukes"],
        col_strats=["Don't use\nnuclear weapons", "Use nuclear\nweapons"],
        payoffs=[[("0", "0"), (None, None)],
                 [(None, None), ("-100", "-100")]],
        caption=True,
        cell_texts={(0, 1): "(Cell empty due to USA's commitment to "
                            "retaliate)",
                    (1, 0): "(Cell empty due to USSR's commitment to "
                            "retaliate)"})
    # red X in the LOWER half of each impossible cell, below the
    # caption text (original layout, Nico 2026-08-09)
    for (r, c) in [(0, 1), (1, 0)]:
        cx, cy = anchors[(r, c, 'cell')]
        hw = Inches(1.0)
        y1, y2 = cy + Inches(0.03), cy + Inches(0.44)
        _add_arrow(slide, (cx - hw, y1), (cx + hw, y2),
                   color=RED, weight_pt=2.0, head=False)
        _add_arrow(slide, (cx - hw, y2), (cx + hw, y1),
                   color=RED, weight_pt=2.0, head=False)
    _add_hierarchical_bullets(
        slide, left=MARGIN, top=Inches(6.0), width=Inches(12.3),
        height=Inches(1.0),
        items=[
            ("(Much) more on commitment in your Strategy class!", 0),
            ("Crucial: Is the player's “commitment” credible?", 1),
        ],
        size=22, sub_size=20, line_spacing_pts=4)
    _draw_footer(slide, FOOTER_TEXT, 68)
    return slide


def slide_69_penalty_kicks(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_GAMES)
    _draw_action_title(
        slide, "Penalty Kicks: Find the (Pure-Strategy) Nash "
               "Equilibrium")
    _add_hierarchical_bullets(
        slide, left=MARGIN, top=Inches(1.65), width=Inches(12.4),
        height=Inches(1.7),
        items=[
            ("Note:", 0),
            ([("Pure Strategy", {'underline': True}),
              (": Each player picks either “Left” or “Right” (as in "
               "all games we've seen thus far)", {})], 1),
            ([("Mixed Strategy", {'underline': True}),
              (": Players can also assign probabilities with which "
               "they pick Left or Right (not relevant for our class)",
               {})], 1),
        ],
        size=22, sub_size=20, line_spacing_pts=6)
    _add_payoff_matrix(
        slide, left=Inches(4.75), top=Inches(4.2), cell_w=Inches(2.3),
        cell_h=Inches(0.95),
        row_player="Kicker", col_player="Goalie",
        row_strats=["Left", "Right"], col_strats=["Left", "Right"],
        payoffs=[[("-1", "1"), ("1", "-1")],
                 [("1", "-1"), ("-1", "1")]],
        caption=True)
    _poll_badge(slide)
    _draw_footer(slide, FOOTER_TEXT, 69)
    return slide


def slide_71_penalty_solution(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_GAMES)
    _draw_action_title(slide, "Penalty Kicks: Solution")
    anchors = _add_payoff_matrix(
        slide, left=Inches(4.75), top=Inches(2.85), cell_w=Inches(2.3),
        cell_h=Inches(1.0),
        row_player="Kicker", col_player="Goalie",
        row_strats=["Left", "Right"], col_strats=["Left", "Right"],
        payoffs=[[("-1", "1"), ("1", "-1")],
                 [("1", "-1"), ("-1", "1")]],
        caption=True)
    # kicker (rows, blue): prefers the side the goalie does NOT cover
    _br_arrow(slide, anchors[(0, 0, 'row')], anchors[(1, 0, 'row')],
              ROW_BLUE)
    _br_circle(slide, anchors[(1, 0, 'row')], ROW_BLUE)
    _br_arrow(slide, anchors[(1, 1, 'row')], anchors[(0, 1, 'row')],
              ROW_BLUE)
    _br_circle(slide, anchors[(0, 1, 'row')], ROW_BLUE)
    # goalie (cols, gold): prefers to match the kicker
    _br_arrow(slide, anchors[(0, 1, 'col')], anchors[(0, 0, 'col')],
              GOLD)
    _br_circle(slide, anchors[(0, 0, 'col')], GOLD)
    _br_arrow(slide, anchors[(1, 0, 'col')], anchors[(1, 1, 'col')],
              GOLD)
    _br_circle(slide, anchors[(1, 1, 'col')], GOLD)
    _add_rounded_filled_box(
        slide, Inches(3.4), Inches(6.1), Inches(6.5), Inches(0.75),
        "There is no (pure-strategy) Nash Equilibrium",
        fill=GOLD, text_color=NAVY, size=20, bold=True,
        corner_pct=0.15)
    _draw_footer(slide, FOOTER_TEXT, 71)
    _set_notes(slide, (
        "This is a nice illustration of what “equilibrium” means: "
        "given that one player goes “right”, the other one also wants "
        "to go “right” – but that's not the case here. Goalie and "
        "Kicker never want to be in the same cell."))
    return slide


def slide_72_advertising_poll(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_GAMES)
    _draw_action_title(slide,
                       "Find the (Pure-Strategy) Nash Equilibrium")
    _add_convention_box(
        slide, Inches(0.35), Inches(5.75), Inches(3.3), Inches(1.1),
        runs=[("Interpretation of payoffs:", {'size': 16,
                                              'underline': True, 'color': NAVY}),
              ("Firms' profits", {'size': 16, 'color': NAVY,
                                  'newline': True})],
        size=16)
    _add_payoff_matrix(
        slide, left=Inches(5.3), top=Inches(2.85), cell_w=Inches(2.75),
        cell_h=Inches(1.05),
        row_player="Coke", col_player="Pepsi",
        row_strats=["Advertise", "Don't\nAdvertise"],
        col_strats=["Advertise", "Don't Advertise"],
        payoffs=[[("150", "150"), ("450", "-75")],
                 [("-75", "450"), ("225", "225")]],
        caption=True)
    _poll_badge(slide)
    _draw_footer(slide, FOOTER_TEXT, 72)
    _set_notes(slide, (
        "The answer is 1. Both advertise, and make lower profits than "
        "if they could coordinate on not advertising. This is just "
        "like the prisoner's dilemma."))
    return slide


def slide_74_pepsi_coke_solution(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_GAMES)
    _draw_action_title(slide, "Pepsi vs. Coke: Solution")
    anchors = _add_payoff_matrix(
        slide, left=Inches(5.3), top=Inches(2.85), cell_w=Inches(2.75),
        cell_h=Inches(1.05),
        row_player="Coke", col_player="Pepsi",
        row_strats=["Advertise", "Don't\nAdvertise"],
        col_strats=["Advertise", "Don't Advertise"],
        payoffs=[[("150", "150"), ("450", "-75")],
                 [("-75", "450"), ("225", "225")]],
        caption=True, nash_cells=[(0, 0)])
    for c in (0, 1):
        _br_arrow(slide, anchors[(1, c, 'row')], anchors[(0, c, 'row')],
                  ROW_BLUE)
        _br_circle(slide, anchors[(0, c, 'row')], ROW_BLUE)
    for r in (0, 1):
        _br_arrow(slide, anchors[(r, 1, 'col')], anchors[(r, 0, 'col')],
                  GOLD)
        _br_circle(slide, anchors[(r, 0, 'col')], GOLD)
    _add_convention_box(
        slide, Inches(0.35), Inches(5.6), Inches(3.3), Inches(1.1),
        runs=[("Interpretation of payoffs:", {'size': 16,
                                              'underline': True, 'color': NAVY}),
              ("Firms' profits", {'size': 16, 'color': NAVY,
                                  'newline': True})],
        size=16)
    _add_rounded_filled_box(
        slide, Inches(4.0), Inches(6.15), Inches(6.6), Inches(0.7),
        "Both advertise — lower profits than if they could coordinate",
        fill=GOLD, text_color=NAVY, size=18, bold=True,
        corner_pct=0.15)
    _draw_footer(slide, FOOTER_TEXT, 74)
    _set_notes(slide, (
        "The answer is 1. Both advertise, and make lower profits than "
        "if they could coordinate on not advertising. This is just "
        "like the prisoner's dilemma."))
    return slide


def slide_76_split_or_steal(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_GAMES)
    _draw_action_title(slide, "“Split or Steal”: Payoff Matrix")
    _add_hierarchical_bullets(
        slide, left=MARGIN, top=Inches(1.65), width=Inches(12.4),
        height=Inches(1.3),
        items=[
            ("We assume that good reputation (when splitting) is worth "
             "£5k (while bad reputation is worth £0)", 0),
            ("We round the total jackpot to £100k", 0),
        ],
        size=22, sub_size=20, line_spacing_pts=6)
    anchors = _add_payoff_matrix(
        slide, left=Inches(4.85), top=Inches(3.55), cell_w=Inches(2.5),
        cell_h=Inches(1.05),
        row_player="Steven", col_player="Sarah",
        row_strats=["Split", "Steal"], col_strats=["Split", "Steal"],
        payoffs=[[("£55", "£55"), ("£5", "£100")],
                 [("£100", "£5"), ("£0", "£0")]],
        caption=True, payoff_size=22)
    # ORIGINAL technique (Nico 2026-08-09): the 4 payoffs hide under
    # gradient cover squares; arrows and the 2 Nash ovals sit ON TOP
    for r in (0, 1):
        for c in (0, 1):
            _add_cover_box(slide,
                           Inches(4.85 + c * 2.5 + 0.04),
                           Inches(3.55 + r * 1.05 + 0.04),
                           Inches(2.42), Inches(0.97))
    _br_arrow(slide, anchors[(0, 0, 'row')], anchors[(1, 0, 'row')],
              ROW_BLUE)
    _br_arrow(slide, anchors[(1, 1, 'row')], anchors[(0, 1, 'row')],
              ROW_BLUE)
    _br_arrow(slide, anchors[(0, 0, 'col')], anchors[(0, 1, 'col')],
              GOLD)
    _br_arrow(slide, anchors[(1, 1, 'col')], anchors[(1, 0, 'col')],
              GOLD)
    for (r, c) in [(0, 1), (1, 0)]:
        ow, oh = int(Inches(2.5 + 0.22)), int(Inches(1.05 + 0.20))
        cxx, cyy = anchors[(r, c, 'cell')]
        oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, cxx - ow // 2,
                                      cyy - oh // 2, ow, oh)
        oval.fill.background()
        oval.line.color.rgb = GOLD
        oval.line.width = Pt(2.75)
        oval.shadow.inherit = False
    _add_rounded_filled_box(
        slide, Inches(10.75), Inches(4.35), Inches(2.35), Inches(0.7),
        "2 Nash Equilibria", fill=GOLD, text_color=NAVY, size=16,
        bold=True, corner_pct=0.25)
    _draw_footer(slide, FOOTER_TEXT, 76)
    return slide


def slide_80_steven_commits(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_GAMES)
    _draw_action_title(slide, "Commitment and Strategy")
    _add_hierarchical_bullets(
        slide, left=MARGIN, top=Inches(1.65), width=Inches(12.4),
        height=Inches(0.9),
        items=[
            ([("If Steven could ", {}),
              ("commit", {'underline': True}),
              (" to Steal, what would be Sarah's optimal response?",
               {})], 0),
        ],
        size=24, sub_size=22, line_spacing_pts=6)
    anchors = _add_payoff_matrix(
        slide, left=Inches(4.85), top=Inches(3.3), cell_w=Inches(2.5),
        cell_h=Inches(1.05),
        row_player="Steven", col_player="Sarah",
        row_strats=["Split", "Steal"], col_strats=["Split", "Steal"],
        payoffs=[[(None, None), (None, None)],
                 [("£100", "£5"), ("£0", "£0")]],
        caption=True, nash_cells=[(1, 0)], payoff_size=22,
        cell_texts={(0, 0): "(Cell impossible due to Steven's "
                            "commitment)",
                    (0, 1): "(Cell impossible due to Steven's "
                            "commitment)"})
    # red X in the LOWER half of each impossible cell, below the
    # caption text (original layout, Nico 2026-08-09)
    for (r, c) in [(0, 0), (0, 1)]:
        cx, cy = anchors[(r, c, 'cell')]
        hw = Inches(1.0)
        y1, y2 = cy + Inches(0.03), cy + Inches(0.46)
        _add_arrow(slide, (cx - hw, y1), (cx + hw, y2),
                   color=RED, weight_pt=2.0, head=False)
        _add_arrow(slide, (cx - hw, y2), (cx + hw, y1),
                   color=RED, weight_pt=2.0, head=False)
    # small X through Steven's 'Split' row label (in the original)
    _add_arrow(slide, (Inches(4.2), Inches(3.71)),
               (Inches(4.76), Inches(3.99)), color=RED,
               weight_pt=2.0, head=False)
    _add_arrow(slide, (Inches(4.2), Inches(3.99)),
               (Inches(4.76), Inches(3.71)), color=RED,
               weight_pt=2.0, head=False)
    _br_arrow(slide, anchors[(1, 1, 'col')], anchors[(1, 0, 'col')],
              GOLD)
    _br_circle(slide, anchors[(1, 0, 'col')], GOLD)
    _add_hierarchical_bullets(
        slide, left=MARGIN, top=Inches(6.05), width=Inches(12.3),
        height=Inches(1.0),
        items=[
            ("(Much) more on this in your Strategy class!", 0),
            ("Crucial: Is the player's “commitment” credible?", 1),
        ],
        size=22, sub_size=20, line_spacing_pts=4)
    _draw_footer(slide, FOOTER_TEXT, 80)
    slide._element.set('show', '0')
    return slide


def slide_81_nick_ibrahim(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_GAMES)
    _draw_action_title(slide, "“Split or Steal”: Payoff Matrix")
    _add_hierarchical_bullets(
        slide, left=MARGIN, top=Inches(1.6), width=Inches(12.4),
        height=Inches(1.7),
        items=[
            ("Nick credibly argues that he will steal. Let's assume he "
             "gets “satisfaction” = 10 from stealing and seriously "
             "plans to share 50-50", 0, {'size': 20}),
            ("Good reputation for splitting is still worth 5", 0,
             {'size': 20}),
            ([("Ibrahim believes that Nick will share winnings with "
               "probability ", {}),
              ("p", {'italic': True, 'color': GOLD}),
              (" > 0", {})], 0, {'size': 20}),
        ],
        size=20, sub_size=18, line_spacing_pts=6)
    _add_payoff_matrix(
        slide, left=Inches(4.85), top=Inches(4.25), cell_w=Inches(2.6),
        cell_h=Inches(1.0),
        row_player="Nick", col_player="Ibrahim",
        row_strats=["Split", "Steal"], col_strats=["Split", "Steal"],
        payoffs=[[("£55", "£55"), ("£5", "£100")],
                 [("£60*", "p×£50"), ("£10", "£0")]],
        caption=True, payoff_size=20)
    _add_text(slide, MARGIN, Inches(6.72), Inches(9.0), Inches(0.35),
              "*Assuming that payoff is shared after stealing. Thus, "
              "payoff = 0.5×100 + 10",
              size=13, italic=True, color=GRAY, font="Calibri")
    _draw_footer(slide, FOOTER_TEXT, 81)
    _set_notes(slide, (
        "Video that will be shown in Strategy: "
        "http://www.youtube.com/watch?v=S0qjK3TWZE8"))
    slide._element.set('show', '0')
    return slide


def slide_82_entry_game(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_GAMES)
    _draw_action_title(slide, "Boeing vs. Airbus: Entry Game")
    _add_hierarchical_bullets(
        slide, left=MARGIN, top=Inches(1.95), width=Inches(7.9),
        height=Inches(4.4),
        items=[
            ("2 firms: Boeing and Airbus", 0),
            ("One type of product: Ultra-long-range (ULR) aircraft", 0),
            ("Decision: produce or not produce", 0),
            ("Boeing: incumbent (777-200LR)", 1),
            ("Airbus: must decide whether to enter (A350 ULR)", 1),
            ("Analyze the situation with and without subsidies", 0),
        ],
        size=26, sub_size=24, line_spacing_pts=14)
    _add_media_image(slide, "image58.png", left=Inches(8.85),
                     top=Inches(2.05), width=Inches(3.9),
                     rounded=True, shadow=True)   # Boeing 777 (Emirates)
    _add_media_image(slide, "image57.png", left=Inches(8.85),
                     top=Inches(4.35), width=Inches(3.9),
                     rounded=True, shadow=True)   # Airbus A350 (Qatar)
    _draw_footer(slide, FOOTER_TEXT, 82)
    _set_notes(slide, (
        "Some background on the 777-200LR, in case it comes up: "
        "https://aeronauticsonline.com/why-it-failed-boeing-777-200lr/"))
    return slide


def _entry_matrix(slide, payoffs, nash_cells, *, top=Inches(3.1)):
    return _add_payoff_matrix(
        slide, left=Inches(5.0), top=top, cell_w=Inches(2.5),
        cell_h=Inches(1.05),
        row_player="Boeing", col_player="Airbus",
        row_strats=["Produce", "Don't\nProduce"],
        col_strats=["Produce", "Don't Produce"],
        payoffs=payoffs, caption=True, nash_cells=nash_cells)


def slide_83_no_subsidy(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_GAMES)
    _draw_action_title(slide, "No Subsidy Case")
    _add_convention_box(
        slide, Inches(0.4), Inches(2.3), Inches(3.55), Inches(1.5),
        runs=[("Interpretation of payoffs:", {'size': 16,
                                              'underline': True, 'color': NAVY}),
              ("Annual profits from ULR aircraft (in $B)",
               {'size': 16, 'color': NAVY, 'newline': True})],
        size=16)
    anchors = _entry_matrix(
        slide, [[("-1", "-1"), ("5", "0")], [("0", "5"), ("0", "0")]],
        [(0, 1)])
    # the second (symmetric) equilibrium gets a DASHED oval
    cx, cy = anchors[(1, 0, 'cell')]
    ow, oh = int(Inches(2.5 + 0.22)), int(Inches(1.05 + 0.20))
    oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - ow // 2,
                                  cy - oh // 2, ow, oh)
    oval.fill.background()
    oval.line.color.rgb = GOLD
    oval.line.width = Pt(2.25)
    oval.shadow.inherit = False
    _dash_shape_line(oval, 'dash')
    for c in (0, 1):
        best_r = 1 if c == 0 else 0
        _br_arrow(slide, anchors[(1 - best_r, c, 'row')],
                  anchors[(best_r, c, 'row')], ROW_BLUE)
        _br_circle(slide, anchors[(best_r, c, 'row')], ROW_BLUE)
    for r in (0, 1):
        best_c = 1 if r == 0 else 0
        _br_arrow(slide, anchors[(r, 1 - best_c, 'col')],
                  anchors[(r, best_c, 'col')], GOLD)
        _br_circle(slide, anchors[(r, best_c, 'col')], GOLD)
    _add_outlined_box(
        slide, Inches(9.3), Inches(5.9), Inches(3.6), Inches(1.05),
        "Boeing is already incumbent. So we're initially in this Nash "
        "equilibrium",
        line=GOLD, text_color=NAVY, size=14, bold=True, line_w=1.75,
        rounded=True, shadow=True, corner_pct=0.15)
    cx2, cy2 = anchors[(0, 1, 'cell')]
    _add_arrow(slide, (Inches(10.9), Inches(5.88)),
               (cx2 + Inches(0.6), cy2 + Inches(0.35)), color=GOLD,
               weight_pt=1.5, head=True)
    _draw_footer(slide, FOOTER_TEXT, 83)
    return slide


def slide_84_subsidy(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_GAMES)
    _draw_action_title(slide, "Subsidy Case")
    _add_hierarchical_bullets(
        slide, left=MARGIN, top=Inches(1.7), width=Inches(12.3),
        height=Inches(0.6),
        items=[("EU pays $2B to Airbus if they produce", 0)],
        size=24, sub_size=22, line_spacing_pts=6)
    anchors = _entry_matrix(
        slide, [[("-5", "1"), ("5", "0")], [("0", "7"), ("0", "0")]],
        [(1, 0)], top=Inches(3.3))
    for c in (0, 1):
        best_r = 1 if c == 0 else 0
        _br_arrow(slide, anchors[(1 - best_r, c, 'row')],
                  anchors[(best_r, c, 'row')], ROW_BLUE)
        _br_circle(slide, anchors[(best_r, c, 'row')], ROW_BLUE)
    for r in (0, 1):
        _br_arrow(slide, anchors[(r, 1, 'col')], anchors[(r, 0, 'col')],
                  GOLD)
        _br_circle(slide, anchors[(r, 0, 'col')], GOLD)
    _add_outlined_box(
        slide, Inches(9.55), Inches(5.95), Inches(3.4), Inches(0.9),
        "Unique Nash equilibrium in the EU-only subsidy case",
        line=GOLD, text_color=NAVY, size=14, bold=True, line_w=1.75,
        rounded=True, shadow=True, corner_pct=0.15)
    cx, cy = anchors[(1, 0, 'cell')]
    _add_arrow(slide, (Inches(10.4), Inches(5.93)),
               (cx + Inches(0.6), cy + Inches(0.4)), color=GOLD,
               weight_pt=1.5, head=True)
    _draw_footer(slide, FOOTER_TEXT, 84)
    return slide


def slide_85_both_subsidies(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_GAMES)
    _draw_action_title(slide, "Subsidy in Both EU and US")
    _add_hierarchical_bullets(
        slide, left=MARGIN, top=Inches(1.7), width=Inches(12.3),
        height=Inches(0.6),
        items=[("US now also pays $2B to Boeing if they produce", 0)],
        size=24, sub_size=22, line_spacing_pts=6)
    anchors = _entry_matrix(
        slide, [[("1", "1"), ("7", "0")], [("0", "7"), ("0", "0")]],
        [(0, 0)], top=Inches(3.3))
    for c in (0, 1):
        _br_arrow(slide, anchors[(1, c, 'row')], anchors[(0, c, 'row')],
                  ROW_BLUE)
        _br_circle(slide, anchors[(0, c, 'row')], ROW_BLUE)
    for r in (0, 1):
        _br_arrow(slide, anchors[(r, 1, 'col')], anchors[(r, 0, 'col')],
                  GOLD)
        _br_circle(slide, anchors[(r, 0, 'col')], GOLD)
    _add_outlined_box(
        slide, Inches(9.55), Inches(5.95), Inches(3.4), Inches(0.9),
        "Unique Nash equilibrium in the EU+US subsidy case",
        line=GOLD, text_color=NAVY, size=14, bold=True, line_w=1.75,
        rounded=True, shadow=True, corner_pct=0.15)
    cx, cy = anchors[(0, 0, 'cell')]
    _add_arrow(slide, (Inches(10.4), Inches(5.93)),
               (cx + Inches(0.6), cy + Inches(0.4)), color=GOLD,
               weight_pt=1.5, head=True)
    _draw_footer(slide, FOOTER_TEXT, 85)
    return slide


# --------------------------------------------------------------------------
# Batch F — recap + commitment backup (85–87)
# --------------------------------------------------------------------------

def slide_86_recap(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_RECAP)
    _draw_action_title(slide, "Recap of Module 7 (see Video on BL)")

    col_w = Inches(6.1)
    left1, left2 = MARGIN, MARGIN + col_w + Inches(0.4)
    _add_rounded_filled_box(slide, left1, Inches(1.6), col_w,
                            Inches(0.55), "Oligopoly", fill=NAVY,
                            text_color=WHITE, size=22, bold=True,
                            corner_pct=0.15)
    _add_rounded_filled_box(slide, left2, Inches(1.6), col_w,
                            Inches(0.55), "Game Theory", fill=NAVY,
                            text_color=WHITE, size=22, bold=True,
                            corner_pct=0.15)
    _add_hierarchical_bullets(
        slide, left=left1 + Inches(0.05), top=Inches(2.35),
        width=col_w - Inches(0.1), height=Inches(4.6),
        items=[
            ("Competition with homogenous products:", 0),
            ("Cournot: Firms compete on quantity. Positive profits", 1),
            ("Equilibrium = intersection of the reaction functions", 2,
             {'size': 14}),
            ("Bertrand: Firms compete on price. Zero profit if both "
             "have same MC", 1),
            ("Price wars → same outcome as perfect competition (if "
             "same MC)", 2, {'size': 14}),
            ("Competition with differentiated products:", 0),
            ("The two products are imperfect substitutes "
             "(differentiated)", 1),
            ("Bertrand: Each firm sets price for its differentiated "
             "product", 1),
            ("If one firm sets a higher price, the other firm also "
             "reacts by raising its price", 2, {'size': 14}),
            ("Equilibrium = intersection of the reaction functions", 2,
             {'size': 14}),
        ],
        size=18, sub_size=16, line_spacing_pts=5)
    _add_hierarchical_bullets(
        slide, left=left2 + Inches(0.05), top=Inches(2.35),
        width=col_w - Inches(0.1), height=Inches(4.6),
        items=[
            ("Strategic thinking: Anticipate other player's reaction", 0),
            ("Strategy: Optimized reaction to each action by other "
             "player", 1),
            ("Use payoff tables", 0),
            ("Dominant Strategy", 1),
            ("Nash Equilibrium", 1),
        ],
        size=20, sub_size=18, line_spacing_pts=8)
    _draw_footer(slide, FOOTER_TEXT, 86)
    return slide


def slide_87_tickets(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_COMMIT)
    _draw_action_title(slide, "Commitment in Selling Tickets")
    _add_media_image(slide, "image59.png", left=Inches(1.05),
                     top=Inches(1.6), width=Inches(11.2),
                     rounded=False, shadow=True)
    _draw_footer(slide, FOOTER_TEXT, 87)
    return slide


def slide_88_tickets_2(prs):
    slide = _blank_slide(prs)
    _draw_top_bar_tc(slide, TAG_COMMIT)
    _draw_action_title(slide, "Commitment in Selling Ticket Sales")
    _add_media_image(slide, "image60.png", left=Inches(2.57),
                     top=Inches(1.42), width=Inches(8.2),
                     rounded=False, shadow=True)
    _add_text(slide, Inches(1.35), Inches(6.28), Inches(10.6),
              Inches(0.5),
              "No sign of falling prices as the game comes closer. "
              "Why?", size=24, bold=True, color=NAVY, font="Calibri",
              align=PP_ALIGN.CENTER)
    # deck-standard back pill (returns to the last viewed slide)
    back = _add_rounded_filled_box(
        slide, Inches(11.72), Inches(6.6), Inches(1.55), Inches(0.46),
        "← Back", fill=NAVY, text_color=WHITE, size=16, bold=True,
        corner_pct=0.5)
    rPr_run = back.text_frame.paragraphs[0].runs[0]
    hl = rPr_run._r.get_or_add_rPr()
    link = hl.makeelement(qn('a:hlinkClick'),
                          {qn('r:id'): '',
                           'action': 'ppaction://hlinkshowjump'
                                     '?jump=lastslideviewed'})
    hl.append(link)
    _draw_footer(slide, FOOTER_TEXT, 88)
    return slide


# --------------------------------------------------------------------------
# Deck assembly — 87 slides, 1:1 with the source deck.
# --------------------------------------------------------------------------

STUB = "TO BUILD — next batch"
STUB_POLL = "PHASE 3 — PollEverywhere slide, spliced from source deck"
STUB_VIDEO = "PHASE 3 — video slide, spliced from source deck"
STUB_EXAMPLE = "TO BUILD — awaiting guiding-example decision (nickel)"


def _group_shapes(slide, shapes, name="HandGroup"):
    """Merge the given shapes into ONE <p:grpSp> (off/ext = bounding
    box, chOff/chExt = off/ext) — build-time implementation of Nico's
    hand groupings (2026-08-12). Members keep absolute positions; the
    group sits at the FIRST member's z-position."""
    from pptx.oxml import parse_xml
    from pptx.oxml.ns import nsdecls
    if len(shapes) < 2:
        return None
    spTree = slide.shapes._spTree
    x0 = min(sh.left for sh in shapes)
    y0 = min(sh.top for sh in shapes)
    x1 = max(sh.left + sh.width for sh in shapes)
    y1 = max(sh.top + sh.height for sh in shapes)
    gid = max(int(e.get('id'))
              for e in spTree.iter(qn('p:cNvPr'))) + 1
    grp = parse_xml(
        '<p:grpSp %s><p:nvGrpSpPr>'
        '<p:cNvPr id="%d" name="%s %d"/>'
        '<p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>'
        '<p:grpSpPr><a:xfrm>'
        '<a:off x="%d" y="%d"/><a:ext cx="%d" cy="%d"/>'
        '<a:chOff x="%d" y="%d"/><a:chExt cx="%d" cy="%d"/>'
        '</a:xfrm></p:grpSpPr></p:grpSp>'
        % (nsdecls('p', 'a'), gid, name, gid,
           x0, y0, x1 - x0, y1 - y0, x0, y0, x1 - x0, y1 - y0))
    spTree.insert(spTree.index(shapes[0]._element), grp)
    for sh in shapes:
        grp.append(sh._element)
    return grp


def _new_shapes_since(slide, before_ids):
    return [sh for sh in slide.shapes
            if sh.shape_id not in before_ids]


def _shape_ids(slide):
    return {sh.shape_id for sh in slide.shapes}


def _group_pic_captions(prs):
    """Group every picture with its caption/source line (Teaching
    CLAUDE.md picture+caption grouping rule, 2026-08-08). A caption is
    a small all-italic (≤13 pt) text box within ~0.3" below (or just
    above) the picture and horizontally centered on it."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.oxml import parse_xml
    from pptx.oxml.ns import nsdecls
    n_groups = 0
    for slide_idx, slide in enumerate(prs.slides, 1):
        spTree = slide.shapes._spTree
        pics = [sh for sh in slide.shapes
                if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
        if not pics:
            continue
        caps = []
        for sh in slide.shapes:
            if sh.shape_type != MSO_SHAPE_TYPE.TEXT_BOX:
                continue
            runs = [r for para in sh.text_frame.paragraphs
                    for r in para.runs if (r.text or "").strip()]
            if runs and all(r.font.italic and r.font.size is not None
                            and r.font.size.pt <= 13 for r in runs):
                caps.append(sh)
        used = set()
        for pic in pics:
            for cap in caps:
                if cap.shape_id in used:
                    continue
                below = cap.top - (pic.top + pic.height)
                above = pic.top - (cap.top + cap.height)
                near = (-Inches(0.05) < below < Inches(0.3)
                        or -Inches(0.05) < above < Inches(0.3))
                ccx = cap.left + cap.width // 2
                pcx = pic.left + pic.width // 2
                centered = abs(ccx - pcx) <= Inches(0.35)
                if not (near and centered):
                    continue
                used.add(cap.shape_id)
                x0 = min(pic.left, cap.left)
                y0 = min(pic.top, cap.top)
                x1 = max(pic.left + pic.width, cap.left + cap.width)
                y1 = max(pic.top + pic.height, cap.top + cap.height)
                gid = max(int(e.get('id')) for e in
                          spTree.iter(qn('p:cNvPr'))) + 1
                grp = parse_xml(
                    '<p:grpSp %s><p:nvGrpSpPr>'
                    '<p:cNvPr id="%d" name="PicCaption %d"/>'
                    '<p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>'
                    '<p:grpSpPr><a:xfrm>'
                    '<a:off x="%d" y="%d"/><a:ext cx="%d" cy="%d"/>'
                    '<a:chOff x="%d" y="%d"/><a:chExt cx="%d" cy="%d"/>'
                    '</a:xfrm></p:grpSpPr></p:grpSp>'
                    % (nsdecls('p', 'a'), gid, gid,
                       x0, y0, x1 - x0, y1 - y0,
                       x0, y0, x1 - x0, y1 - y0))
                spTree.insert(spTree.index(pic._element), grp)
                grp.append(pic._element)
                grp.append(cap._element)
                n_groups += 1
                break
    print(f"picture+caption groups: {n_groups}")


def _add_ticket_jump(prs):
    """Slide-67 jump link to the ticket-commitment example on slide 86
    (the ORIGINAL deck's bottom-right action button), rendered as a
    navy pill in the standard corner position."""
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    slides = list(prs.slides)
    s67, s86 = slides[67], slides[86]   # displays 68 and 87 after the
    #                                     concept-map insert (2026-08-10)
    box = _add_rounded_filled_box(
        s67, Inches(11.72), Inches(6.6), Inches(1.55), Inches(0.46),
        "Tickets  →", fill=NAVY, text_color=WHITE, size=14, bold=True,
        corner_pct=0.5, shadow=True)
    rid = s67.part.relate_to(s86.part, RT.SLIDE)
    cNvPr = box._element.find(qn('p:nvSpPr') + '/' + qn('p:cNvPr'))
    hl = ET.SubElement(cNvPr, qn('a:hlinkClick'))
    hl.set(qn('r:id'), rid)
    hl.set('action', 'ppaction://hlinksldjump')


def build(out_path=None):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_title(prs)                                            # 1
    slide_02_recap(prs)                                            # 2
    slide_03_roadmap(prs)                                          # 3
    make_module_outline(prs, 4, descriptions=True)                 # 4
    slide_05_concept_map(prs)                                      # 5
    slide_06_market_structures(prs)                                # 6
    slide_07_spectrum(prs)                                         # 7
    slide_08_characteristics(prs)                                  # 8
    slide_09_concentration(prs)                                    # 9
    slide_10_nickel(prs)                                           # 10
    make_module_outline(prs, 11, section_tag=TAG_COLLU,
                        part_idx=0, sub_idx=0)                     # 11
    slide_12_collusion(prs)                                        # 12
    slide_13_cartels(prs)                                          # 13
    slide_14_adm(prs)                                              # 14
    make_stub(prs, 15, TAG_COLLU, "The ADM Price-Fixing Scandal",
              "copy video over by hand (not done yet to keep "
              "file size manageable)")   # 14 — Nico's note 2026-08-10
    slide_16_adm_chart(prs)                                        # 16
    make_module_outline(prs, 17, section_tag=TAG_COURNOT,
                        part_idx=0, sub_idx=1)                     # 17
    slide_18_note(prs)                                             # 18
    slide_19_two_models(prs)                                       # 19
    slide_20_two_models_bullets(prs)                               # 20
    slide_21_duopoly_setup(prs)                                    # 21
    slide_22_cournot_assumptions(prs)                              # 22
    slide_23_reaction_i(prs)                                       # 23
    slide_24_reaction_ii(prs)                                      # 24
    slide_25_reaction_function(prs)                                # 25
    slide_26_cournot_equilibrium(prs)                              # 26
    slide_27_computation(prs)                                      # 27
    slide_28_further_examples(prs)                                 # 28
    make_module_outline(prs, 29, section_tag=TAG_BERTRAND,
                        part_idx=0, sub_idx=2)                     # 29
    slide_30_bertrand(prs)                                         # 30
    slide_31_bertrand_assumptions(prs)                             # 31
    slide_32_concrete_poll(prs)                                    # 32
    make_stub(prs, 33, TAG_BERTRAND, "Poll", STUB_POLL)            # 33
    slide_34_bertrand_chart(prs)                                   # 34
    slide_35_outcomes_comparison(prs)                              # 35
    slide_36_outcomes_math(prs)                                    # 36
    slide_37_overview(prs)                                         # 37
    make_module_outline(prs, 38, section_tag=TAG_DIFF,
                        part_idx=0, sub_idx=3)                     # 38
    slide_39_diff_products(prs)                                    # 39
    slide_40_diff_assumptions(prs)                                 # 40
    slide_41_diff_setup(prs)                                       # 41
    slide_42_airbus_i(prs)                                         # 42
    slide_43_airbus_ii(prs)                                        # 43
    slide_44_equilibrium_prices(prs)                               # 44
    slide_45_bruinlearn(prs)                                       # 45
    slide_46_takeaways(prs)                                        # 46
    slide_47_concluding_discussion(prs)                            # 47
    make_module_outline(prs, 48, section_tag=TAG_GT,
                        part_idx=1, sub_idx=0)                     # 48
    slide_49_what_is_game_theory(prs)                              # 49
    slide_50_warmup_game(prs)                                      # 50
    slide_51_beauty_contest(prs)                                   # 51
    slide_52_numerical_beauty(prs)                                 # 52
    make_stub(prs, 53, TAG_GT, "Poll", STUB_POLL)                  # 53
    slide_54_key_concepts(prs)                                     # 54
    make_module_outline(prs, 65, section_tag=TAG_GAMES,
                        part_idx=1, sub_idx=1)                     # 55
    slide_56_pd_setup(prs)                                         # 56
    slide_57_pd_p1(prs)                                            # 57
    slide_58_pd_p2(prs)                                            # 58
    slide_59_pd_equilibrium(prs)                                   # 59
    slide_60_solve_dilemma(prs)                                    # 60
    slide_61_podcast(prs)                                          # 61
    slide_62_fair_play(prs)                                        # 62
    slide_63_trade_wars(prs)                                       # 63
    slide_64_collective_action(prs)                                # 64
    slide_65_chicken(prs)                                          # 65
    slide_66_mad(prs)                                              # 66
    slide_67_group_work(prs)                                       # 67
    slide_68_commitment(prs)                                       # 68
    slide_69_penalty_kicks(prs)                                    # 69
    make_stub(prs, 70, TAG_GAMES, "Poll", STUB_POLL)               # 70
    slide_71_penalty_solution(prs)                                 # 71
    slide_72_advertising_poll(prs)                                 # 72
    make_stub(prs, 73, TAG_GAMES, "Poll", STUB_POLL)               # 73
    slide_74_pepsi_coke_solution(prs)                              # 74
    make_stub(prs, 75, TAG_GAMES, "Game Theory on TV: Split or Steal?",
              STUB_VIDEO)                                          # 75
    slide_76_split_or_steal(prs)                                   # 76
    make_stub(prs, 77, TAG_GAMES, "Split or Steal: Coordination?",
              STUB_VIDEO)                                          # 77
    make_stub(prs, 78, TAG_GAMES, "Poll", STUB_POLL)               # 78
    make_stub(prs, 79, TAG_GAMES, "Split or Steal", STUB_VIDEO)    # 79
    slide_80_steven_commits(prs)                                   # 80
    slide_81_nick_ibrahim(prs)                                     # 81
    slide_82_entry_game(prs)                                       # 82
    slide_83_no_subsidy(prs)                                       # 83
    slide_84_subsidy(prs)                                          # 84
    slide_85_both_subsidies(prs)                                   # 85
    slide_86_recap(prs)                                            # 86
    slide_87_tickets(prs)                                          # 87
    slide_88_tickets_2(prs)                                        # 88

    _add_ticket_jump(prs)
    _group_pic_captions(prs)
    out = Path(out_path) if out_path else OUT_DIR / "Module 7 - Revised.pptx"
    prs.save(str(out))
    print(f"saved {out} — {len(prs.slides._sldIdLst)} slides")
    return out


if __name__ == "__main__":
    import sys
    build(sys.argv[1] if len(sys.argv) > 1 else None)
