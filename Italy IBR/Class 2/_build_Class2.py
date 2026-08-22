"""
Build "Class 2 - Revised.pptx" -- Italy IBR (MGMTEX 421), Class 2.

Adapts the 2023 4:3 deck `Class 2.pptx` to the Class 1 visual language
(13.33 x 7.5", navy/gold/gray palette, single master, live slide-number
fields).  Chrome geometry is lifted verbatim from `Class 1 - Revised.pptx`.

Pipeline (each pass is rerunnable, phase-3 style):
    python _build_Class2.py          -> builds the 15 script-buildable slides
    python _splice_poll.py           -> inserts the PollEverywhere slide at #4
    python _animate.py all apply     -> injects the <p:timing> builds

The PollEverywhere slide cannot go through python-pptx (it strips the poll
`tags` relationship), so it is spliced by OOXML surgery afterwards.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.oxml.ns import qn
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from lxml import etree

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(HERE, "Class 2 - Revised.pptx")
IMG = os.path.join(HERE, "_source_images")

# ----------------------------------------------------------------- palette
NAVY = "0B2B4E"
GOLD = "E09F3E"
GRAY = "555B66"
RULE = "C8CDD3"
CREAM = "FDF6E6"
WHITE = "FFFFFF"

# ------------------------------------------------- canvas / chrome geometry
SLIDE_W, SLIDE_H = 12192000, 6858000
LEFT = 251999                 # standard left margin
CONTENT_W = 11687697          # standard content width
BAR_W, BAR_H = 12191695, 384048

FOOTER_TEXT = "International Business Residential – Italy"

NSDECL = (
    'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
    'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" '
    'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
)

_STYLE = (
    '<p:style><a:lnRef idx="1"><a:schemeClr val="accent1"/></a:lnRef>'
    '<a:fillRef idx="3"><a:schemeClr val="accent1"/></a:fillRef>'
    '<a:effectRef idx="2"><a:schemeClr val="accent1"/></a:effectRef>'
    '<a:fontRef idx="minor"><a:schemeClr val="lt1"/></a:fontRef></p:style>'
)
SHADOW = ('<a:effectLst><a:outerShdw blurRad="50800" dist="38100" dir="2700000" '
          'rotWithShape="0"><a:srgbClr val="000000"><a:alpha val="30000"/>'
          '</a:srgbClr></a:outerShdw></a:effectLst>')
SHADOW_SOFT = ('<a:effectLst><a:outerShdw blurRad="40000" dist="25400" dir="2700000" '
               'rotWithShape="0"><a:srgbClr val="000000"><a:alpha val="22000"/>'
               '</a:srgbClr></a:outerShdw></a:effectLst>')


def esc(t):
    return (t.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;"))


def E(inches):
    """Inches -> integer EMU (decimal EMUs break PowerPoint silently)."""
    return int(round(inches * 914400))


# ------------------------------------------------------------- XML builders
_ids = {"n": 100}


def nid():
    _ids["n"] += 1
    return _ids["n"]


def run(text, sz, b=0, i=0, color=NAVY, u=None, link=None):
    """One <a:r> run spec."""
    return dict(text=text, sz=sz, b=b, i=i, color=color, u=u, link=link)


def _run_xml(r):
    rpr = '<a:rPr lang="en-US" sz="%d" b="%d" i="%d"' % (r["sz"], r["b"], r["i"])
    rpr += ' u="sng"' if r.get("u") else ""
    rpr += ">"
    if r.get("link"):
        rpr += ('<a:hlinkClick r:id="%s" action="ppaction://hlinksldjump"/>'
                % r["link"])
    else:
        # schema order: fill BEFORE latin, else PowerPoint ignores the color
        rpr += '<a:solidFill><a:srgbClr val="%s"/></a:solidFill>' % r["color"]
        rpr += '<a:latin typeface="Calibri"/>'
    rpr += "</a:rPr>"
    return "<a:r>%s<a:t>%s</a:t></a:r>" % (rpr, esc(r["text"]))


def para(runs, algn="l", lvl=None, bullet=None, spc_bef=None, line=None):
    return dict(runs=runs, algn=algn, lvl=lvl, bullet=bullet,
                spc_bef=spc_bef, line=line)


def _para_xml(p):
    ppr = "<a:pPr"
    if p.get("bullet") == 1:
        ppr += ' marL="342900" indent="-274320"'
    elif p.get("bullet") == 2:
        ppr += ' marL="731520" lvl="1" indent="-274320"'
    if p.get("algn"):
        ppr += ' algn="%s"' % p["algn"]
    ppr += ">"
    if p.get("line"):
        ppr += '<a:lnSpc><a:spcPct val="%d"/></a:lnSpc>' % p["line"]
    if p.get("spc_bef") is not None:
        ppr += '<a:spcBef><a:spcPts val="%d"/></a:spcBef>' % p["spc_bef"]
    if p.get("bullet"):
        ch = "▪" if p["bullet"] == 1 else "–"
        ppr += ('<a:buClr><a:srgbClr val="%s"/></a:buClr>'
                '<a:buFont typeface="Calibri"/><a:buChar char="%s"/>' % (NAVY, ch))
    ppr += "</a:pPr>"
    return "<a:p>%s%s</a:p>" % (ppr, "".join(_run_xml(r) for r in p["runs"]))


def textbox(name, x, y, w, h, paras, anchor=None, autofit=True, wrap="square"):
    body = '<a:bodyPr wrap="%s" lIns="0" tIns="0" rIns="0" bIns="0"' % wrap
    if anchor:
        body += ' anchor="%s"' % anchor
    body += "><a:spAutoFit/></a:bodyPr>" if autofit else "/>"
    return (
        '<p:sp><p:nvSpPr><p:cNvPr id="%d" name="%s"/><p:cNvSpPr txBox="1"/>'
        '<p:nvPr/></p:nvSpPr><p:spPr><a:xfrm><a:off x="%d" y="%d"/>'
        '<a:ext cx="%d" cy="%d"/></a:xfrm><a:prstGeom prst="rect"><a:avLst/>'
        '</a:prstGeom><a:noFill/></p:spPr><p:txBody>%s<a:lstStyle/>%s</p:txBody>'
        "</p:sp>" % (nid(), name, x, y, w, h, body,
                     "".join(_para_xml(p) for p in paras))
    )


def rect(name, x, y, w, h, fill):
    return (
        '<p:sp><p:nvSpPr><p:cNvPr id="%d" name="%s"/><p:cNvSpPr/><p:nvPr/>'
        '</p:nvSpPr><p:spPr><a:xfrm><a:off x="%d" y="%d"/><a:ext cx="%d" cy="%d"/>'
        '</a:xfrm><a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
        '<a:solidFill><a:srgbClr val="%s"/></a:solidFill><a:ln><a:noFill/></a:ln>'
        "<a:effectLst/></p:spPr>%s<p:txBody><a:bodyPr rtlCol=\"0\" anchor=\"ctr\"/>"
        '<a:lstStyle/><a:p><a:pPr algn="ctr"/><a:endParaRPr/></a:p></p:txBody></p:sp>'
        % (nid(), name, x, y, w, h, fill, _STYLE)
    )


def roundrect(name, x, y, w, h, fill, adj=12000, line=None, lw=12700,
              paras=None, shadow=True, anchor="ctr", wrap="square",
              lins=45720, link=None, soft=False):
    ln = ('<a:ln w="%d"><a:solidFill><a:srgbClr val="%s"/></a:solidFill></a:ln>'
          % (lw, line)) if line else "<a:ln><a:noFill/></a:ln>"
    eff = (SHADOW_SOFT if soft else SHADOW) if shadow else ""
    hl = ('<a:hlinkClick r:id="%s" action="ppaction://hlinksldjump"/>' % link
          ) if link else ""
    body = ('<a:bodyPr wrap="%s" lIns="%d" tIns="0" rIns="%d" bIns="0" anchor="%s"/>'
            % (wrap, lins, lins, anchor))
    txt = "".join(_para_xml(p) for p in (paras or [para([], algn="ctr")]))
    return (
        '<p:sp><p:nvSpPr><p:cNvPr id="%d" name="%s">%s</p:cNvPr><p:cNvSpPr/>'
        '<p:nvPr/></p:nvSpPr><p:spPr><a:xfrm><a:off x="%d" y="%d"/>'
        '<a:ext cx="%d" cy="%d"/></a:xfrm><a:prstGeom prst="roundRect"><a:avLst>'
        '<a:gd name="adj" fmla="val %d"/></a:avLst></a:prstGeom>'
        '<a:solidFill><a:srgbClr val="%s"/></a:solidFill>%s%s</p:spPr>'
        "<p:txBody>%s<a:lstStyle/>%s</p:txBody></p:sp>"
        % (nid(), name, hl, x, y, w, h, adj, fill, ln, eff, body, txt)
    )


def ellipse(name, x, y, w, h, fill, line=None, lw=22225):
    ln = ('<a:ln w="%d"><a:solidFill><a:srgbClr val="%s"/></a:solidFill></a:ln>'
          % (lw, line)) if line else "<a:ln><a:noFill/></a:ln>"
    return (
        '<p:sp><p:nvSpPr><p:cNvPr id="%d" name="%s"/><p:cNvSpPr/><p:nvPr/>'
        '</p:nvSpPr><p:spPr><a:xfrm><a:off x="%d" y="%d"/><a:ext cx="%d" cy="%d"/>'
        '</a:xfrm><a:prstGeom prst="ellipse"><a:avLst/></a:prstGeom>'
        '<a:solidFill><a:srgbClr val="%s"/></a:solidFill>%s</p:spPr>'
        '<p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:endParaRPr/></a:p></p:txBody>'
        "</p:sp>" % (nid(), name, x, y, w, h, fill, ln)
    )


def parallelogram(name, x, y, w, h, text, sz=2000):
    return (
        '<p:sp><p:nvSpPr><p:cNvPr id="%d" name="%s"/><p:cNvSpPr/><p:nvPr/>'
        '</p:nvSpPr><p:spPr><a:xfrm><a:off x="%d" y="%d"/><a:ext cx="%d" cy="%d"/>'
        '</a:xfrm><a:prstGeom prst="parallelogram"><a:avLst>'
        '<a:gd name="adj" fmla="val 30000"/></a:avLst></a:prstGeom>'
        '<a:solidFill><a:srgbClr val="%s"/></a:solidFill><a:ln><a:noFill/></a:ln>'
        '<a:effectLst><a:outerShdw blurRad="50800" dist="38100" dir="2700000" '
        'rotWithShape="0"><a:srgbClr val="000000"><a:alpha val="28000"/></a:srgbClr>'
        "</a:outerShdw></a:effectLst></p:spPr><p:txBody>"
        '<a:bodyPr wrap="none" lIns="0" rIns="0" anchor="ctr"/><a:lstStyle/>%s'
        "</p:txBody></p:sp>"
        % (nid(), name, x, y, w, h, GOLD,
           _para_xml(para([run(text, sz, b=1, color=NAVY)], algn="ctr")))
    )


# ------------------------------------------------------------------- chrome
def top_bar(tag):
    return [
        rect("Rectangle 1", 0, 0, BAR_W, BAR_H, NAVY),
        textbox("TextBox 2", LEFT, 0, 10972800, BAR_H,
                [para([run("Italy IBR  ·  " + tag, 1600, b=1, color=WHITE)])],
                anchor="ctr"),
    ]


def title_block(title, sz=3000, y=566928, h=640080):
    return [
        textbox("TextBox 3", LEFT, y, CONTENT_W, h,
                [para([run(title, sz, b=1, color=NAVY)])]),
        rect("Rectangle 4", LEFT, 1188720, CONTENT_W, 18288, RULE),
        rect("Rectangle 5", LEFT, 1175004, 2011680, 45720, GOLD),
    ]


def footer(page=True, idx=1):
    out = [
        rect("Rectangle 7", 0, 6537960, BAR_W, 18288, RULE),
        rect("Rectangle 8", LEFT, 6524244, 2011680, 45720, GOLD),
    ]
    if page:
        out.append(textbox("TextBox 9", LEFT, 6583680, 10058400, 292608,
                           [para([run(FOOTER_TEXT, 1200, color=GRAY)])]))
        guid = "{5B7A6A11-0000-4A00-9C00-%012d}" % idx
        out.append(
            '<p:sp><p:nvSpPr><p:cNvPr id="%d" name="TextBox 10"/>'
            '<p:cNvSpPr txBox="1"/><p:nvPr/></p:nvSpPr><p:spPr><a:xfrm>'
            '<a:off x="11430000" y="6583680"/><a:ext cx="548640" cy="292608"/>'
            '</a:xfrm><a:prstGeom prst="rect"><a:avLst/></a:prstGeom><a:noFill/>'
            '</p:spPr><p:txBody><a:bodyPr wrap="square" lIns="0" tIns="0" '
            'rIns="0" bIns="0"><a:spAutoFit/></a:bodyPr><a:lstStyle/><a:p>'
            '<a:pPr algn="r"/><a:fld id="%s" type="slidenum"><a:rPr lang="en-US" '
            'sz="1200"><a:solidFill><a:srgbClr val="%s"/></a:solidFill>'
            '<a:latin typeface="Calibri"/></a:rPr><a:t>%d</a:t></a:fld>'
            '<a:endParaRPr lang="en-US" sz="1200"><a:solidFill>'
            '<a:srgbClr val="%s"/></a:solidFill><a:latin typeface="Calibri"/>'
            "</a:endParaRPr></a:p></p:txBody></p:sp>"
            % (nid(), guid, GRAY, idx, GRAY))
    return out


def back_button(rid, x=E(11.72), y=E(6.6)):
    return roundrect("BackButton", x, y, E(1.55), E(0.46), NAVY, adj=30000,
                     shadow=True, link=rid,
                     paras=[para([run("← Back", 1400, b=1, color=WHITE)],
                                 algn="ctr")])


def pill(name, x, y, w, h, text, sz=1400, fill=NAVY, color=WHITE, adj=30000):
    return roundrect(name, x, y, w, h, fill, adj=adj, shadow=True, wrap="none",
                     lins=0, paras=[para([run(text, sz, b=1, color=color)],
                                         algn="ctr")])


# --------------------------------------------------------------- deck data
COMPANIES = [
    dict(logo="_logo_kering.png", commons="File:Kering-logo.svg", group="Group 2",
         title="Kering – A Portfolio of Iconic Luxury Houses",
         sector="Luxury Group", slot="Wednesday 9 Sep  ·  09.00 – 10.30",
         blurb=("A session with a senior leader at one of the world’s most "
                "powerful luxury conglomerates, exploring how Kering builds and "
                "manages a portfolio of iconic houses."),
         logo_h=1.55),
    dict(logo="_logo_prada.png", commons="File:Prada Group - logo (Italy).svg", group="Group 5",
         title="Prada Group – Global Strategy in Luxury Fashion",
         sector="Luxury Fashion", slot="Tuesday 8 Sep  ·  15.00 – 16.30",
         blurb=("Visit to the company’s headquarters for a session with a "
                "senior leader to explore how a global luxury brand approaches "
                "global strategy."),
         logo_h=0.80),
    dict(logo="_logo_inter.png", commons="File:FC Internazionale Milano 2021.svg", group="Group 6",
         title="Inter Milan FC – A Football Club as a Global Business",
         sector="Iconic Sports Team", slot="Wednesday 9 Sep  ·  15.00 – 16.30",
         blurb=("Explore how Inter operates as a global sports organization and "
                "how its brand and commercial strategy drives success on and off "
                "the pitch."),
         logo_h=2.35),
    dict(logo="_logo_pirelli.png", commons="File:Pirelli - logo full (Italy, 1997).svg", group="Group 8",
         title="Pirelli – R&D and Innovation Behind the Tire",
         sector="Vehicle Manufacturing", slot="Tuesday 8 Sep  ·  09.00 – 10.45",
         blurb=("A session with the CTO of a global leader in tire manufacturing, "
                "exploring how R&D and innovation drive commercial strategy across "
                "motorsport and luxury markets."),
         logo_h=1.20),
    dict(logo="_logo_essilor.png", commons="File:Logo EssilorLuxottica.svg", group="Group 9",
         title="EssilorLuxottica – Premium Eyewear at Global Scale",
         sector="Eyewear Manufacturer", slot="Tuesday 8 Sep  ·  11.30 – 13.00",
         blurb=("Hear from a global leader in premium eyewear design and "
                "manufacturing, to learn how the company partners with top "
                "fashion brands."),
         logo_h=0.62),
]

# Each debate lists its two sides in ASCENDING GROUP-NUMBER order, so the lower
# group always takes the left column and reveals first.
DEBATES = [
    dict(title="Italy and the EU / the Euro",
         sides=[("CONTRA", "Group 1",
                 "Italy would be economically better off outside the EU, "
                 "without the Euro"),
                ("PRO", "Group 3",
                 "Being in the EU and using the Euro is good for Italy from an "
                 "economic standpoint")]),
    # PRO/CONTRA here map to Team B/Team A of the menu file, not A/B: Group 4
    # argues the optimistic side (confirmed by Nico 2026-08-22).
    dict(title="Italy’s Future in Tech and AI",
         sides=[("PRO", "Group 4",
                 "Italy is already on a good path"),
                ("CONTRA", "Group 7",
                 "Italy is lagging far behind on tech and AI")]),
]

# Running order = by group number. A debate sits at its LOWER group number, so
# the two sides of one question stay on one slide.
#   Groups 1 & 3 -> EU/Euro debate | 2 Kering | 4 & 7 tech-AI debate
#   5 Prada | 6 Inter | 8 Pirelli | 9 EssilorLuxottica
SEQUENCE = [
    ("debate", 0), ("company", 0), ("debate", 1),
    ("company", 1), ("company", 2), ("company", 3), ("company", 4),
]

LINEUP = [
    ("Italy and the EU / the Euro", "Debate  ·  Groups 1 & 3"),
    ("Kering", "Company  ·  Group 2"),
    ("Italy’s Future in Tech and AI", "Debate  ·  Groups 4 & 7"),
    ("Prada Group", "Company  ·  Group 5"),
    ("Inter Milan FC", "Company  ·  Group 6"),
    ("Pirelli", "Company  ·  Group 8"),
    ("EssilorLuxottica", "Company  ·  Group 9"),
]

ROADMAP = [
    "Logistics – the schedule from here",
    "Arrival in Milan – a quick poll",
    "Seven group presentations – in order of group number",
]

NOTES = {}


# ------------------------------------------------------------------ engine
def add_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])   # 6 = Blank


def put(slide, xmls):
    tree = slide.shapes._spTree
    for x in xmls:
        frag = ('<root %s>%s</root>' % (NSDECL, x))
        root = etree.fromstring(frag.encode("utf8"))
        for child in list(root):
            tree.append(child)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def picture(slide, path, x, y, w=None, h=None, flat=True):
    pic = slide.shapes.add_picture(path, Emu(x), Emu(y),
                                   Emu(w) if w else None, Emu(h) if h else None)
    if not flat:
        spPr = pic._element.spPr
        spPr.append(etree.fromstring(
            '<a:prstGeom %s prst="roundRect"><a:avLst>'
            '<a:gd name="adj" fmla="val 4500"/></a:avLst></a:prstGeom>'
            % NSDECL))
        spPr.append(etree.fromstring(
            '<a:effectLst %s><a:outerShdw blurRad="50800" dist="38100" '
            'dir="2700000" rotWithShape="0"><a:srgbClr val="000000">'
            '<a:alpha val="32000"/></a:srgbClr></a:outerShdw></a:effectLst>'
            % NSDECL))
    return pic


def link_to(slide, target_slide):
    return slide.part.relate_to(target_slide.part, RT.SLIDE)


# ------------------------------------------------------------------- build
def main():
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H

    s = {}
    for k in range(15):
        s[k] = add_slide(prs)

    # ------------------------------------------------------------ 1  title
    put(s[0], [
        textbox("TextBox 1", LEFT, E(2.35), CONTENT_W, E(1.30),
                [para([run("Presentations and Debates", 6000, b=1, color=NAVY)],
                      algn="ctr")]),
        textbox("TextBox 2", LEFT, E(3.70), CONTENT_W, E(0.75),
                [para([run("Class 2  ·  Topics on the Italian Economy",
                           3200, b=1, color=GOLD)], algn="ctr")]),
        rect("Rectangle 3", E(4.665), E(4.75), E(4.00), E(0.06), GOLD),
        textbox("TextBox 4", LEFT, E(5.15), CONTENT_W, E(0.55),
                [para([run("MGMTEX 421  ·  International Business "
                           "Residential in Italy", 2400, b=1, color=GRAY)],
                      algn="ctr")]),
        textbox("TextBox 5", LEFT, E(5.85), CONTENT_W, E(0.50),
                [para([run("Prof. Nico Voigtländer  ·  UCLA Anderson",
                           2000, color=GRAY)], algn="ctr")]),
    ] + footer(page=False))
    notes(s[0], "Welcome back. This is our second class, and today the class is "
                "yours: seven groups present, and we debate. Two broad-topic "
                "debates and five of the companies we will visit in Milan and "
                "Turin. I will keep time and moderate, and we close by looking "
                "ahead to the trip itself.")

    # --------------------------------------------------------- 2  roadmap
    put(s[1], top_bar("Today’s Roadmap")
        + title_block("Class Today: Your Presentations and Debates")
        + footer(idx=2))
    pitch, bh = E(0.86), E(0.62)
    y0 = E(4.22) - (len(ROADMAP) * pitch - (pitch - bh)) // 2
    rm = []
    for i, t in enumerate(ROADMAP):
        y = y0 + i * pitch
        cy = y + (bh - E(0.50)) // 2
        rm.append(roundrect("RmBand", LEFT, y, CONTENT_W, bh, CREAM, adj=16000,
                            line=GOLD, lw=19050, soft=True))
        rm.append(ellipse("RmOval", E(0.60), cy, E(0.50), E(0.50), NAVY, line=GOLD))
        rm.append(textbox("RmNum", E(0.60), cy, E(0.50), E(0.50),
                          [para([run(str(i + 1), 1700, b=1, color=WHITE)],
                                algn="ctr")], anchor="ctr", autofit=False))
        rm.append(textbox("RmTxt", E(1.34), y, E(11.6), bh,
                          [para([run(t, 2400, b=1, color=NAVY)])],
                          anchor="ctr", autofit=False))
    put(s[1], rm)
    notes(s[1], "Here is the plan. First the logistics from here to Milan, then a "
                "quick poll on when you arrive, and then the main event: the "
                "seven group presentations. Two of them are head-to-head debates "
                "on broad questions about the Italian economy, and five are "
                "company presentations on firms we will actually walk into "
                "during the trip.")

    # ------------------------------------------------------- 3  logistics
    put(s[2], top_bar("Course Logistics")
        + title_block("Logistics – The Schedule from Here")
        + footer(idx=3))
    sched_rid = link_to(s[2], s[14])
    LB = [
        para([run("Class today: discussion of the presentation debates", 2800)],
             bullet=1),
        para([run("Present a 2 – 3 min executive summary", 2400)],
             bullet=2, spc_bef=300),
        para([run("Followed by discussion", 2400)], bullet=2, spc_bef=300),
        para([run("Max. 20 minutes per topic", 2400)], bullet=2, spc_bef=300),
        para([run("Trip to Italy: Sep 6 – 12 – Milan & Turin   (", 2800),
              run("Schedule", 2800, u=True, link=sched_rid),
              run(")", 2800)], bullet=1, spc_bef=1200),
        para([run("Note: casual dress code", 2400)], bullet=2, spc_bef=300),
        para([run("Get melatonin for the jetlag!", 2400)], bullet=2, spc_bef=300),
        para([run("Class on Oct 4: debrief and discussion of student reports",
                  2800)], bullet=1, spc_bef=1200),
    ]
    put(s[2], [textbox("TextBox 6", LEFT, E(1.75), CONTENT_W, E(4.90), LB,
                       anchor="ctr", autofit=False)])
    notes(s[2], "Three things on timing. Today each group gives a 2 to 3 minute "
                "executive summary of the video you already submitted, and then "
                "we open it up; I will cap each topic at 20 minutes so everyone "
                "gets the floor. The trip runs September 6 to 12, Milan first and "
                "then Turin, and the full day-by-day schedule is on the last "
                "slide if you click through. Dress code is casual throughout, and "
                "do get melatonin before you fly. We meet once more on October 4 "
                "to debrief and talk through your reports.")

    # -------------------------------------------------------- 4  divider
    put(s[3], top_bar("Presentations and Debates") + [
        textbox("TextBox 3", LEFT, E(3.00), CONTENT_W, E(1.00),
                [para([run("Presentations and Debates", 4400, b=1, color=NAVY)],
                      algn="ctr")]),
        rect("Rectangle 4", E(5.165), E(4.20), E(3.00), E(0.05), GOLD),
    ] + footer(idx=5))
    notes(s[3], "Now over to you. Seven presentations, in the order on the next "
                "slide.")

    # -------------------------------------------------------- 5  line-up
    put(s[4], top_bar("Presentations and Debates")
        + title_block("Today’s Line-Up – Seven Presentations")
        + footer(idx=6))
    y0, pitch, bh = E(1.86), E(0.70), E(0.58)
    lu = []
    for i, (topic, who) in enumerate(LINEUP):
        y = y0 + i * pitch
        cy = y + (bh - E(0.44)) // 2
        lu.append(roundrect("LuBand", LEFT, y, CONTENT_W, bh, CREAM, adj=16000,
                            line=GOLD, lw=19050, soft=True))
        lu.append(ellipse("LuOval", E(0.60), cy, E(0.44), E(0.44), NAVY, line=GOLD))
        lu.append(textbox("LuNum", E(0.60), cy, E(0.44), E(0.44),
                          [para([run(str(i + 1), 1500, b=1, color=WHITE)],
                                algn="ctr")], anchor="ctr", autofit=False))
        lu.append(textbox("LuTxt", E(1.28), y, E(7.6), bh,
                          [para([run(topic, 2200, b=1, color=NAVY)])],
                          anchor="ctr", autofit=False))
        lu.append(textbox("LuWho", E(8.95), y, E(4.05), bh,
                          [para([run(who, 1800, b=1, color=GRAY)], algn="r")],
                          anchor="ctr", autofit=False))
    put(s[4], lu)
    notes(s[4], "This is the running order, and it follows the sign-up sheet. We "
                "start with the two debates, where two groups argue opposite "
                "sides of the same question, and then move to the five company "
                "presentations. Each company group introduces the firm and then "
                "makes the case for or against investing in it.")

    # ------------------------------------------- 7-13  the presentations
    # One loop over SEQUENCE, so the deck follows group order.
    def emit_debate(sl, d, idx):
        put(sl, top_bar("Broad-Topic Debates")
            + title_block(d["title"])
            + footer(idx=idx))
        colw, gap = E(6.20), E(0.38)
        hdr_y, card_y, card_h = E(1.72), E(2.40), E(3.55)
        shapes = []
        for col, (side, group, claim) in enumerate(d["sides"]):
            x = LEFT + col * (colw + gap)
            shapes.append(
                roundrect("Hdr%d" % (col + 1), x, hdr_y, colw, E(0.56), NAVY,
                          adj=22000,
                          paras=[para([run("%s  ·  %s" % (side, group), 2000,
                                           b=1, color=WHITE)], algn="ctr")]))
            shapes.append(
                roundrect("Card%d" % (col + 1), x, card_y, colw, card_h, CREAM,
                          adj=6000, line=GOLD, lw=15875, lins=E(0.28),
                          paras=[para([run("“" + claim + "”", 2600, i=1,
                                           color=NAVY)], algn="ctr", line=95000)]))
        shapes.append(
            roundrect("Takeaway", LEFT, E(6.30), CONTENT_W, E(0.52), GOLD,
                      adj=18000,
                      paras=[para([run("2 – 3 min executive summary from each "
                                       "side, then open debate", 1900, b=1, i=1,
                                       color=NAVY)], algn="ctr")]))
        put(sl, shapes)
        first, second = d["sides"][0], d["sides"][1]
        notes(sl, "This is a head-to-head debate: the two groups were handed "
                  "opposite sides of the same question, so treat the statements "
                  "on the slide as positions to argue, not as my view. %s opens "
                  "with a 2 to 3 minute executive summary, %s follows, and then "
                  "we open the floor. Listen for which side leans on evidence "
                  "and which leans on assertion."
              % (first[1], second[1]))

    def emit_company(sl, c, idx):
        put(sl, top_bar("Company Presentations")
            + title_block(c["title"])
            + footer(idx=idx))
        # left panel: logo, sector, group pill
        lp_cx = E(3.05)                       # centre of the left column
        lh = E(c["logo_h"])
        from PIL import Image
        im = Image.open(os.path.join(IMG, c["logo"]))
        lw_ = int(lh * im.width / im.height)
        maxw, maxh = E(4.60), E(2.28)   # maxh keeps tall logos off the sector line
        if lw_ > maxw:
            lw_, lh = maxw, int(maxw * im.height / im.width)
        if lh > maxh:
            lw_, lh = int(maxh * im.width / im.height), maxh
        picture(sl, os.path.join(IMG, c["logo"]),
                lp_cx - lw_ // 2, E(2.85) - lh // 2, lw_, lh, flat=True)
        put(sl, [
            textbox("Sector", lp_cx - E(2.30), E(4.18), E(4.60), E(0.40),
                    [para([run(c["sector"], 2000, b=1, color=NAVY)], algn="ctr")]),
            pill("GroupPill", lp_cx - E(1.05), E(4.80), E(2.10), E(0.52),
                 c["group"], sz=1800),
        ])
        # right panel: brief + visit cards
        rx, rw = E(5.90), E(7.15)
        put(sl, [
            roundrect("BriefCard", rx, E(1.72), rw, E(1.90), CREAM, adj=8000,
                      line=NAVY, lw=12700, anchor="t", lins=E(0.26),
                      paras=[
                          para([run("Your brief", 2000, b=1, color=NAVY)],
                               algn="l", spc_bef=1400),
                          para([run("Introduce the company", 2000, color=NAVY)],
                               bullet=1, spc_bef=700),
                          para([run("Then argue why (not) to invest in it", 2000,
                                    color=NAVY)], bullet=1, spc_bef=400),
                      ]),
            roundrect("VisitCard", rx, E(3.86), rw, E(2.35), CREAM, adj=8000,
                      line=NAVY, lw=12700, anchor="t", lins=E(0.26),
                      paras=[
                          para([run("We visit them on the trip", 2000, b=1,
                                    color=NAVY)], algn="l", spc_bef=1400),
                          para([run(c["slot"], 2000, b=1, color=GOLD)],
                               algn="l", spc_bef=700),
                          para([run(c["blurb"], 1800, i=1, color=NAVY)],
                               algn="l", spc_bef=600, line=95000),
                      ]),
        ])
        notes(sl, "%s is on the trip: we see them %s. The group has 2 to 3 "
                  "minutes to introduce the company and then to argue why we "
                  "should, or should not, invest in it. Keep the investment case "
                  "concrete, and note the questions you want to put to the "
                  "company when we are in the room with them.  Sector: %s. "
                  "[Logo: Wikimedia Commons, %s]"
              % (c["title"].split(" – ")[0],
                 c["slot"].replace("  ·  ", ", "), c["sector"], c["commons"]))

    for pos, (kind, which) in enumerate(SEQUENCE):
        sl = s[5 + pos]
        disp = 7 + pos            # display number once the poll sits at #4
        if kind == "debate":
            emit_debate(sl, DEBATES[which], disp)
        else:
            emit_company(sl, COMPANIES[which], disp)

    # ------------------------------------------------------ 13 thank you
    put(s[12], [
        textbox("TextBox 1", LEFT, E(3.00), CONTENT_W, E(1.30),
                [para([run("Thank You!", 6000, b=1, color=NAVY)], algn="ctr")]),
        rect("Rectangle 2", E(5.165), E(4.30), E(3.00), E(0.06), GOLD),
    ] + footer(page=False))
    notes(s[12], "Thank you all. See you in Milan on September 6, and we will "
                 "regroup on Zoom on October 4 to debrief.")

    # -------------------------------------------------- 14 backup divider
    put(s[13], top_bar("Backup Slides") + [
        textbox("TextBox 3", LEFT, E(3.00), CONTENT_W, E(1.00),
                [para([run("Backup Slides", 4400, b=1, color=NAVY)], algn="ctr")]),
        rect("Rectangle 4", E(5.165), E(4.20), E(3.00), E(0.05), GOLD),
    ] + footer(idx=15))

    # ------------------------------------------------- 15 trip schedule
    put(s[14], top_bar("Backup")
        + title_block("Trip Schedule (Sep 6 – 12)")
        + footer(idx=16))
    ih = E(5.42)
    iw = int(ih * 2434 / 1546)
    ix = (SLIDE_W - iw) // 2
    picture(s[14], os.path.join(IMG, "itinerary_2026.png"),
            ix, E(1.36), iw, ih, flat=True)
    back = link_to(s[14], s[2])
    put(s[14], [
        back_button(back),
        textbox("SchedSrc", ix, E(6.84), iw, E(0.26),
                [para([run("Itinerary: Legacy Ventures, “UCLA Italy 2026”",
                           1200, i=1, color=GRAY)], algn="ctr")]),
    ])
    notes(s[14], "This is the full week. Milan Sunday through Wednesday, then we "
                 "coach over to Turin on Thursday morning. The company sessions "
                 "are the ones five of the groups are presenting on today: "
                 "Pirelli and EssilorLuxottica and Prada on Tuesday, Kering and "
                 "Inter on Wednesday, then Lavazza and Carioca in Turin, and "
                 "Stellantis and LeapFactory on Friday. Bocconi opens the week "
                 "on Monday morning, and the Wednesday late-morning slot is "
                 "still being settled between Dolce&Gabbana and Eataly.")

    prs.save(OUT)
    print("wrote %s - %d slides" % (OUT, len(prs.slides._sldIdLst)))


if __name__ == "__main__":
    main()
