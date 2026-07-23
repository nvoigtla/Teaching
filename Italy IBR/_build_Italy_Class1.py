"""
============================================================================
  STALE  --  DO NOT RE-RUN THIS SCRIPT.  (frozen 2026-07-19)

  "Class 1 - Revised.pptx" has since been HAND-EDITED in PowerPoint and
  patched in place (_apply_edits.py). It now contains changes this script
  does NOT reproduce, including:
    - the TA email on slide 4, transport rewording + NH Milano Touring +
      two pasted hotel photos on slide 6, the pasted schedule image on
      slide 99;
    - the in-place edits: slide-4 "Schedule" -> slide-99 jump, slide-99
      "Back" button, refreshed transport prices, styled hotel photos.
  Re-running would REGENERATE from scratch and OVERWRITE all of the above.

  The .pptx is now the SOURCE OF TRUTH. Make further changes IN PLACE via
  OOXML/lxml surgery (see _apply_edits.py / _splice_polls.py / _wire_links.py),
  never by re-running this file. Kept for its helpers and as the build record.
============================================================================

Build the reformatted "Class 1 - Revised.pptx" (Italy IBR, MGMTEX 421) in the
405 clean visual language.

SCOPE: slides 1-20 of the original 123-slide "Class 1.pptx".  This revised
deck is 16:9 and grows as later subsets are reformatted.  The original 4:3
"Class 1.pptx" stays untouched as the source.

TWO-STEP BUILD (always run in this order):
    1. python _build_Italy_Class1.py        -> builds all script-buildable slides
    2. python _splice_polls.py               -> injects the 3 live PollEverywhere
                                                embeds (slides 11, 13, 20)
Because both steps are scripted, re-running them in sequence is idempotent.

Reuses primitives from the Module 3 template (single source of the 405 style):
    _build_template_samples.py
"""

import sys
from io import BytesIO
from pathlib import Path

from lxml import etree as ET
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

HERE = Path(__file__).parent
TEMPLATE_DIR = HERE.parent / "Module 3"
sys.path.insert(0, str(TEMPLATE_DIR))

from _build_template_samples import (  # noqa: E402
    FADED, GOLD, GOLD_W, GRAY, MARGIN, NAVY, RULE, RULE_W,
    SLIDE_H, SLIDE_W, WHITE,
    _add_rect, _add_text, _blank_slide, _set_bullet_char,
)

IMG = HERE / "Images"
OUT = HERE / "Class 1 - Revised.pptx"

SECTION_PREFIX = "Italy IBR"
FOOTER_TEXT = "International Business Residential – Italy"
SYLLABUS = "Course Syllabus -- IBR Italy Sep 2026.pdf"  # relative link target

# The 9-section roadmap (slide 10 overview + the you-are-here dividers).
ROADMAP = [
    "(Very) Early History of Italy",
    "The Roman Empire",
    "The Dark Ages",
    "Independent City States (Communes)",
    "The Renaissance",
    "Decline and Napoleon",
    "Italian Unification",
    "Fascism and World War II",
    "The Italian Economy Today",
]


# ==========================================================================
# Chrome helpers (title-case top bar + live-field footer)
# ==========================================================================

def _tag(section):
    return f"{SECTION_PREFIX}  ·  {section}"


def _top_bar(slide, section):
    bar_h = Inches(0.42)
    _add_rect(slide, 0, 0, SLIDE_W, bar_h, NAVY)
    _add_text(slide, MARGIN, 0, Inches(12), bar_h, _tag(section),
              size=16, bold=True, color=WHITE, font="Calibri",
              anchor=MSO_ANCHOR.MIDDLE)


def _action_title(slide, title):
    _add_text(slide, MARGIN, Inches(0.62), RULE_W, Inches(0.7),
              title, size=30, bold=True, color=NAVY, font="Calibri")
    _add_rect(slide, MARGIN, Inches(1.30), RULE_W, Inches(0.02), RULE)
    _add_rect(slide, MARGIN, Inches(1.285), GOLD_W, Inches(0.05), GOLD)


_FLD_GUID = "{{5B7A6A11-0000-4A00-9C00-00000000{:04d}}}"


def _footer(slide, page_num):
    """Footer rule + gold accent + course text + LIVE slidenum field."""
    _add_rect(slide, 0, Inches(7.15), SLIDE_W, Inches(0.02), RULE)
    _add_rect(slide, MARGIN, Inches(7.135), GOLD_W, Inches(0.05), GOLD)
    _add_text(slide, MARGIN, Inches(7.20), Inches(11), Inches(0.32),
              FOOTER_TEXT, size=12, color=GRAY)
    box = _add_text(slide, Inches(12.5), Inches(7.20), Inches(0.6),
                    Inches(0.32), "", size=12, color=GRAY,
                    align=PP_ALIGN.RIGHT)
    # Replace the empty run with a slidenum field so numbering is live.
    p = box.text_frame.paragraphs[0]._p
    fld = ET.SubElement(p, qn('a:fld'))
    fld.set('id', _FLD_GUID.format(page_num))
    fld.set('type', 'slidenum')
    rPr = ET.SubElement(fld, qn('a:rPr'))
    rPr.set('lang', 'en-US'); rPr.set('sz', '1200')
    fill = ET.SubElement(rPr, qn('a:solidFill'))
    ET.SubElement(fill, qn('a:srgbClr')).set('val', '555B66')
    ET.SubElement(rPr, qn('a:latin')).set('typeface', 'Calibri')
    t = ET.SubElement(fld, qn('a:t')); t.text = str(page_num)


def _set_notes(slide, text):
    tf = slide.notes_slide.notes_text_frame
    tf.clear()
    tf.text = text


# ==========================================================================
# Hierarchical bullets with inline hyperlink runs
# ==========================================================================
# Item form:  (level, content[, style])
#   content : str  OR  list of run-segments
#   segment : (text,)                      plain run
#             (text, {"link": url})        hyperlink run (navy, underlined)
#             (text, {"bold": True})       etc.
# style (optional dict): {"space_before": pts}
#
# Level sizing (course canon): L0 24pt / L1 22pt / L2 20pt.

_LVL_SIZE = {0: 24, 1: 22, 2: 20}
_LVL_MARL = {0: 342900, 1: 731520, 2: 1097280}
_LVL_INDENT = -274320
_LVL_CHAR = {0: "▪", 1: "–", 2: "·"}   # ▪ – ·


def _emit_run(p, text, *, size, color=NAVY, bold=False, italic=False,
              underline=False, link=None):
    run = p.add_run()
    run.text = text
    f = run.font
    f.name = "Calibri"; f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = color
    if underline:
        f.underline = True
    if link:
        run.hyperlink.address = link
    return run


def _hbullets(slide, items, *, left=MARGIN, top=Inches(1.9),
              width=RULE_W, height=Inches(5.0), anchor=MSO_ANCHOR.MIDDLE,
              main_space=16, sub_space=5):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if anchor is not None:
        tf.vertical_anchor = anchor

    for i, item in enumerate(items):
        level, content = item[0], item[1]
        style = item[2] if len(item) > 2 else {}
        size = _LVL_SIZE[level]

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        pPr = p._p.get_or_add_pPr()
        pPr.set('marL', str(_LVL_MARL[level]))
        pPr.set('indent', str(_LVL_INDENT))
        if level > 0:
            pPr.set('lvl', str(level))
        # space-before
        sb = style.get("space_before", main_space if level == 0 else sub_space)
        if i > 0:
            spc = ET.SubElement(pPr, qn('a:spcBef'))
            ET.SubElement(spc, qn('a:spcPts')).set('val', str(sb * 100))

        # runs
        segs = [(content, {})] if isinstance(content, str) else content
        for seg in segs:
            txt = seg[0]
            o = seg[1] if len(seg) > 1 else {}
            _emit_run(p, txt, size=size,
                      color=o.get("color", NAVY),
                      bold=o.get("bold", False),
                      italic=o.get("italic", False),
                      underline=o.get("underline", bool(o.get("link"))),
                      link=o.get("link"))

        # Pass level-specific marL/indent — _set_bullet_char would otherwise
        # reset marL to its level-0 default and flatten the sub-bullet indent.
        _set_bullet_char(p, char=_LVL_CHAR[level], color=NAVY, font="Calibri",
                         mar_l=_LVL_MARL[level], indent=_LVL_INDENT)
    return box


# ==========================================================================
# Picture placement (rounded corners + soft shadow, aspect preserved)
# ==========================================================================

def _fit(path, max_w, max_h):
    w, h = Image.open(path).size
    r = min(max_w / w, max_h / h)
    return int(w * r), int(h * r)


def _place_image(slide, path, *, cx=None, cy=None, left=None, top=None,
                 max_w, max_h, shadow=True, rounded=True):
    w, h = _fit(str(path), int(max_w), int(max_h))
    if left is None:
        left = int(cx - w / 2)
    if top is None:
        top = int(cy - h / 2)
    pic = slide.shapes.add_picture(str(path), left, top, w, h)
    spPr = pic._element.spPr
    if rounded:
        # Modify the geom python-pptx already emitted (a 2nd prstGeom is
        # schema-invalid and makes PowerPoint reject the whole file).
        geom = spPr.find(qn('a:prstGeom'))
        if geom is None:
            geom = ET.SubElement(spPr, qn('a:prstGeom'))
        geom.set('prst', 'roundRect')
        av = geom.find(qn('a:avLst'))
        if av is None:
            av = ET.SubElement(geom, qn('a:avLst'))
        for gd in av.findall(qn('a:gd')):
            av.remove(gd)
        gd = ET.SubElement(av, qn('a:gd'))
        gd.set('name', 'adj'); gd.set('fmla', f'val {int(0.055 * 100000)}')
    if shadow:
        _shadow(spPr)
    return pic, left, top, w, h


def _shadow(spPr):
    # Remove any existing effectLst (e.g. the empty one python-pptx inserts
    # for shadow.inherit=False) — two effectLst elements corrupt the file.
    for el in spPr.findall(qn('a:effectLst')):
        spPr.remove(el)
    lst = ET.SubElement(spPr, qn('a:effectLst'))
    sh = ET.SubElement(lst, qn('a:outerShdw'))
    sh.set('blurRad', '50800'); sh.set('dist', '38100')
    sh.set('dir', '2700000'); sh.set('rotWithShape', '0')
    c = ET.SubElement(sh, qn('a:srgbClr')); c.set('val', '000000')
    ET.SubElement(c, qn('a:alpha')).set('val', '32000')


def _white_backing(slide, left, top, w, h):
    """White rect w/ shadow behind a table/graphicframe (can't shadow those)."""
    r = _add_rect(slide, int(left), int(top), int(w), int(h), WHITE)
    _shadow(r._element.spPr)
    return r


# ==========================================================================
# SLIDES
# ==========================================================================

def s01_cover(prs):
    slide = _blank_slide(prs)
    _place_image(slide, IMG / "s01_image1.png",
                 cx=SLIDE_W // 2, cy=SLIDE_H // 2,
                 max_w=Inches(11.2), max_h=Inches(6.0),
                 shadow=True, rounded=True)


def s02_title(prs):
    slide = _blank_slide(prs)
    _add_text(slide, MARGIN, Inches(2.35), RULE_W, Inches(1.3),
              "The Italian Economy", size=60, bold=True, color=NAVY,
              font="Calibri", align=PP_ALIGN.CENTER)
    _add_text(slide, MARGIN, Inches(3.7), RULE_W, Inches(0.75),
              "Class 1", size=40, bold=True, color=GOLD,
              font="Calibri", align=PP_ALIGN.CENTER)
    aw = Inches(4.0)
    _add_rect(slide, (SLIDE_W - aw) // 2, Inches(4.75), aw, Inches(0.06), GOLD)
    _add_text(slide, MARGIN, Inches(5.15), RULE_W, Inches(0.55),
              "MGMTEX 421  ·  International Business Residential in Italy",
              size=24, bold=True, color=GRAY, font="Calibri",
              align=PP_ALIGN.CENTER)
    _add_text(slide, MARGIN, Inches(5.85), RULE_W, Inches(0.5),
              "Prof. Nico Voigtländer  ·  UCLA Anderson",
              size=20, color=GRAY, font="Calibri", align=PP_ALIGN.CENTER)
    _add_rect(slide, 0, Inches(7.15), SLIDE_W, Inches(0.02), RULE)
    _add_rect(slide, MARGIN, Inches(7.135), GOLD_W, Inches(0.05), GOLD)


def _content(prs, page, section, title, items, notes=None, extras=None,
             bullets_top=Inches(1.9), bullets_h=Inches(5.0)):
    slide = _blank_slide(prs)
    _top_bar(slide, section)
    _action_title(slide, title)
    _hbullets(slide, items, top=bullets_top, height=bullets_h)
    if extras:
        extras(slide)
    _footer(slide, page)
    if notes:
        _set_notes(slide, notes)
    return slide


LOG = "Course Logistics"


def s03(prs):
    _content(prs, 3, LOG, "About the Instructor", [
        (0, "Professor of Global Economics and Management"),
        (1, [("Email: ",), ("nico.v@ucla.edu", {"link": "mailto:nico.v@ucla.edu"})]),
        (0, "Background"),
        (1, "Born and raised in Germany"),
        (1, "Studied Economics and Engineering in Berlin and Cambridge (MA)"),
        (1, "PhD in Economics in Barcelona"),
        (1, "Research on firm productivity, trade, and the role of history in development"),
    ])


def s04(prs):
    _content(prs, 4, LOG, "Logistics – Schedule", [
        (0, "Jul 26 (Zoom): Introduction to the Italian Economy"),
        (1, "Quiz on the assigned cases (individual homework before class)"),
        (0, "Aug 23 (Zoom): Group presentation debates"),
        (1, "Presentation videos (5–10 min) due Aug 17"),
        (0, [("Sep 6–12: Trip to Italy – Milan & Turin   (",),
             ("Schedule", {"link": SYLLABUS}), (")",)]),
        (0, "Oct 4 (Zoom): Debrief and discussion of paper topics"),
        (0, "TA for the class: Elisabetta Campagna"),
    ])


def s05(prs):
    _content(prs, 5, LOG, "Logistics – Transportation", [
        (0, "Getting to Milan – two airports"),
        (1, "Linate (LIN): closer to the center; taxi to hotel 15–20 min"),
        (1, "Malpensa (MXP): direct flights from the U.S., but farther out"),
        (1, [("From Malpensa: ",),
             ("Express train", {"link": "https://www.malpensaexpress.it/en/"}),
             (" 54 min / €13 (every 30 min), ",),
             ("coach", {"link": "https://www.milanomalpensa-airport.com/en/from-to/by-coach"}),
             (" €10, or ",),
             ("taxi", {"link": "https://www.milanomalpensa-airport.com/en/from-to/by-taxi"}),
             (" €104 (“tariffa fissa”)",)]),
        (0, "Return via Turin (TRN): ~15–20 min outside the center"),
    ])


def s06(prs):
    _content(prs, 6, LOG, "Logistics – Local Transport & Hotels", [
        (0, "Local transport"),
        (1, "Walking, metro, bus"),
        (1, "Taxi; Uber (Black only, pricier than cabs)"),
        (0, [("Hotel in Milan: ",),
             ("Starhotel Ritz", {"link": "https://www.starhotels.com/en/our-hotels/ritz-milan/"}),
             (" – near center & central station",)]),
        (0, [("Hotel in Turin: ",),
             ("NH Collection Santo Stefano", {"link": "https://www.nh-hotels.com/hotel/nh-collection-torino-santo-stefano"}),
             (" – right in the center",)]),
    ])


def s07(prs):
    slide = _blank_slide(prs)
    _top_bar(slide, LOG)
    _action_title(slide, "Logistics – Grades")
    rows = [
        ("Quiz – “Italy: The Good, the Bad, and the Ugly”", "10%"),
        ("Video presentations (group basis)", "25%"),
        ("Individual participation", "25%"),
        ("Final individual report", "40%"),
    ]
    tbl_w, tbl_h = Inches(9.2), Inches(3.2)
    tbl_l = (SLIDE_W - tbl_w) // 2
    tbl_t = Inches(2.5)
    _white_backing(slide, tbl_l, tbl_t, tbl_w, tbl_h)
    gf = slide.shapes.add_table(5, 2, tbl_l, tbl_t, tbl_w, tbl_h)
    tb = gf.table
    tb.first_row = False; tb.horz_banding = False
    tb.columns[0].width = Inches(7.4); tb.columns[1].width = Inches(1.8)
    hdr = [("Component", PP_ALIGN.LEFT), ("Weight", PP_ALIGN.CENTER)]
    for c, (txt, al) in enumerate(hdr):
        _fill_cell(tb.cell(0, c), txt, NAVY, WHITE, bold=True, align=al, size=20)
    for r, (comp, wt) in enumerate(rows, start=1):
        bg = WHITE if r % 2 else RGBColor(0xFD, 0xF6, 0xE6)
        _fill_cell(tb.cell(r, 0), comp, bg, NAVY, align=PP_ALIGN.LEFT, size=20)
        _fill_cell(tb.cell(r, 1), wt, bg, NAVY, bold=True,
                   align=PP_ALIGN.CENTER, size=20)
    _footer(slide, 7)


def _fill_cell(cell, text, fill, color, *, bold=False, align=PP_ALIGN.LEFT,
               size=20):
    cell.fill.solid(); cell.fill.fore_color.rgb = fill
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.18); cell.margin_right = Inches(0.12)
    cell.margin_top = Inches(0.04); cell.margin_bottom = Inches(0.04)
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = "Calibri"; r.font.size = Pt(size)
    r.font.bold = bold; r.font.color.rgb = color


def s08(prs):
    _content(prs, 8, LOG, "Logistics – Video Presentations", [
        (0, "Groups of 4–6 students; choose from the topics list"),
        (0, "Prepare a 5–10 min video (e.g., record your group on Zoom)"),
        (1, "Videos due Aug 17 (upload on BruinLearn)"),
        (0, "Present a 2 min executive summary in class Aug 23, then Q&A"),
        (0, [("Sign up in groups: email the TA your topic choices ",),
             ("here", {"link": "https://docs.google.com/forms/d/e/1FAIpQLSdq3S-NP2tCRT4UUd-8kuRC-YgKPzgIJOLj17O27dh6wKLCBQ/viewform"}),
             (" by midnight today",)]),
    ])


def s09(prs):
    _content(prs, 9, LOG, "Logistics – Final Paper", [
        (0, "Due Oct 4 – last Zoom class (6–8 pages)"),
        (0, "Possible topics"),
        (1, "The future of the Italian economy"),
        (1, "The future of the economy in northern Italy"),
        (1, "Other suggestions welcome – need at least 5 papers per topic"),
        (0, "Refer explicitly to what you learned in class and in-country"),
        (0, "Discussed in the last class (may bring tables/graphs only)"),
    ])


def s10_roadmap(prs):
    slide = _blank_slide(prs)
    _top_bar(slide, "Today’s Roadmap")
    _action_title(slide, "Class Today: The Italian Economy (and Its History)")
    items = [(0, s, {"space_before": 8}) for s in ROADMAP]
    _hbullets(slide, items, top=Inches(1.75), height=Inches(5.1),
              main_space=8)
    _footer(slide, 10)


CREAM = RGBColor(0xFD, 0xF6, 0xE6)


def _divider(prs, page, current_idx, subtitle=None):
    """405-style section anchor: the current roadmap section sits in a
    cream/gold 'you are here' band with a filled number badge; the other
    eight sections are dimmed. The slide title is the current section."""
    slide = _blank_slide(prs)
    _top_bar(slide, "Section Divider")
    _action_title(slide, ROADMAP[current_idx])
    _add_text(slide, Inches(9.4), Inches(0.66), Inches(3.5), Inches(0.5),
              f"Section {current_idx + 1} of {len(ROADMAP)}",
              size=16, bold=True, italic=True, color=GOLD,
              font="Calibri", align=PP_ALIGN.RIGHT)

    n = len(ROADMAP)
    row_h = int(Inches(0.56))
    start_y = int(Inches(4.30)) - row_h * n // 2
    badge_d = int(Inches(0.44))
    band_h = int(Inches(0.50))
    bx = int(MARGIN) + int(Inches(0.32))
    text_x = int(MARGIN) + int(Inches(1.02))
    text_w = int(RULE_W) - int(Inches(1.3))

    for i, sec in enumerate(ROADMAP):
        cur = (i == current_idx)
        cy = start_y + row_h * i + row_h // 2
        if cur:
            band = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, int(MARGIN),
                cy - band_h // 2, int(RULE_W), band_h)
            band.fill.solid(); band.fill.fore_color.rgb = CREAM
            band.line.color.rgb = GOLD; band.line.width = Pt(1.5)
            band.shadow.inherit = False
            try:
                band.adjustments[0] = 0.16
            except Exception:
                pass
            _shadow(band._element.spPr)
        # number badge
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, bx, cy - badge_d // 2,
                                       badge_d, badge_d)
        badge.fill.solid(); badge.fill.fore_color.rgb = NAVY if cur else WHITE
        badge.line.color.rgb = GOLD if cur else FADED
        badge.line.width = Pt(1.75 if cur else 1.0)
        badge.shadow.inherit = False
        _add_text(slide, bx, cy - badge_d // 2, badge_d, badge_d, str(i + 1),
                  size=15, bold=True, color=WHITE if cur else FADED,
                  font="Calibri", align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # section name
        _add_text(slide, text_x, cy - row_h // 2, text_w, row_h, sec,
                  size=23 if cur else 21, bold=cur,
                  color=NAVY if cur else FADED, font="Calibri",
                  anchor=MSO_ANCHOR.MIDDLE)
        if cur and subtitle:
            _add_text(slide, int(SLIDE_W - MARGIN) - int(Inches(3.9)),
                      cy - row_h // 2, int(Inches(3.5)), row_h, subtitle,
                      size=16, bold=True, italic=True, color=GOLD,
                      font="Calibri", align=PP_ALIGN.RIGHT,
                      anchor=MSO_ANCHOR.MIDDLE)
    _footer(slide, page)


def s12_divider(prs):
    _divider(prs, 12, 0, subtitle="10,000 BC – 500 BC")


def s17_divider(prs):
    _divider(prs, 17, 1)


def _content_with_images(prs, page, section, title, items, images, notes=None):
    slide = _blank_slide(prs)
    _top_bar(slide, section)
    _action_title(slide, title)
    # bullets on the left ~58% width
    text_w = Inches(6.9)
    _hbullets(slide, items, left=MARGIN, top=Inches(1.9),
              width=text_w, height=Inches(5.0))
    images(slide)
    _footer(slide, page)
    if notes:
        _set_notes(slide, notes)
    return slide


def s14(prs):
    def imgs(slide):
        # map (locator) top, Sassi photo below, right column
        col_cx = Inches(10.55)
        _place_image(slide, IMG / "s14_image5.png", cx=col_cx, cy=Inches(3.0),
                     max_w=Inches(3.7), max_h=Inches(2.35))
        _place_image(slide, IMG / "s14_image4.png", cx=col_cx, cy=Inches(5.45),
                     max_w=Inches(3.7), max_h=Inches(2.15))
    _content_with_images(prs, 14, "Early History",
        "Settlement Before the Roman Empire", [
        (0, "Many different cultures across the peninsula"),
        (0, "City of Matera (Southern Italy)"),
        (1, "Settled since 10,000 BC; famous “Sassi” ancient town"),
        (0, "Rock drawings in Valcamonica (Camuni people)"),
        (1, "~8,000–1,000 BC; world’s largest prehistoric petroglyph collection"),
        (1, "Famous “The Astronauts,” with the Camunian Rose"),
    ], imgs)


def s15(prs):
    def extra(slide):
        _add_text(slide, Inches(6.7), Inches(6.55), Inches(6.35), Inches(0.35),
                  "New Museum on Etruscan history in Milan",
                  size=13, italic=True, color=GRAY, align=PP_ALIGN.RIGHT)
        # attach link to that caption run
        r = slide.shapes[-1].text_frame.paragraphs[0].runs[0]
        r.hyperlink.address = "https://museo.fondazioneluigirovati.org/en"
    _content(prs, 15, "Early History", "Etruscan Civilization", [
        (0, "Central Italy, from around 800 BC"),
        (0, "Chiefdom / tribes → centralized, aristocratic state"),
        (0, "Economy based on mining and metal trade"),
        (0, "Advanced art and a written language"),
        (0, "Expansion led to conflict with the Greeks"),
        (1, "Battles in 600–500 BC weakened the Etruscans"),
        (0, "Assimilated by Rome around 500 BC"),
    ], extras=extra)


def s16(prs):
    def imgs(slide):
        _place_image(slide, IMG / "s16_image6.png", cx=Inches(10.3),
                     cy=Inches(4.35), max_w=Inches(4.3), max_h=Inches(4.9))
    _content_with_images(prs, 16, "Early History",
        "Ethnic Groups in Italy around 400 BC", [
        (0, "Many ethnic groups – Celts in the north, Greeks in the south"),
        (0, "After 500 BC, the Latins of Rome grew in power and influence"),
        (0, "Process of Romanization, with adoption of the Latin language"),
    ], imgs, notes="More info on Celts: https://en.wikipedia.org/wiki/Celts")


ROME = "The Roman Empire"


def s18(prs):
    def imgs(slide):
        _place_image(slide, IMG / "s18_image7.jpeg", cx=Inches(10.4),
                     cy=Inches(4.3), max_w=Inches(4.1), max_h=Inches(4.4))
    _content_with_images(prs, 18, ROME, "Foundation of Rome", [
        (0, "Rome founded in 753 BC by the twin brothers Romulus and Remus"),
        (0, "Legend: they were raised by a she-wolf"),
        (0, "Strategic location"),
        (1, "Ford to cross the river Tiber; easy-to-defend hills"),
        (0, "First a kingdom; from 509 BC a Republic"),
    ], imgs)


def s19(prs):
    _content(prs, 19, ROME, "The Roman Republic", [
        (0, "System of checks and balances"),
        (1, "Magistrates (executive: censor, consul…), elected by the People – plebeians & patricians"),
        (1, "Senate: powerful advisory assembly of unelected aristocrats"),
        (1, "Popular assemblies: voted on legislative, electoral, and judicial matters"),
        (0, [("1st-century-BC crisis → Caesar becomes dictator (49 BC), murdered (44 BC)",)]),
    ])


# ---- Poll slides: chrome + POLL pill + poll snapshot image -----------------
# The live PollEverywhere tags are injected afterward by _splice_polls.py.

def _poll(prs, page, section, question, img_name, notes):
    slide = _blank_slide(prs)
    _top_bar(slide, section)
    _action_title(slide, question)
    _poll_pill(slide)
    _place_image(slide, IMG / img_name, cx=SLIDE_W // 2, cy=Inches(4.35),
                 max_w=Inches(7.4), max_h=Inches(4.9),
                 shadow=True, rounded=False)
    _footer(slide, page)
    _set_notes(slide, notes)
    return slide


def _poll_pill(slide):
    pw, ph = Inches(1.05), Inches(0.34)
    px = SLIDE_W - MARGIN - pw
    py = Inches(0.62)
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px, py, pw, ph)
    shp.fill.solid(); shp.fill.fore_color.rgb = NAVY
    shp.line.fill.background(); shp.shadow.inherit = False
    try:
        shp.adjustments[0] = 0.5
    except Exception:
        pass
    _add_text(slide, px, py, pw, ph, "POLL", size=12, bold=True, color=WHITE,
              font="Calibri", align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    ds = Inches(0.14)
    _add_rect(slide, int(px - Inches(0.16) - ds), int(py + (ph - ds) // 2),
              ds, ds, GOLD)


PN = ("\nPoll Title: Do not modify the notes in this section to avoid "
      "tampering with the Poll Everywhere activity.\nMore info at "
      "polleverywhere.com/support\n\n")


def s11_poll(prs):
    _poll(prs, 11, "Early History", "When did economic activity in Italy begin?",
          "s11_image2.png",
          PN + "When did economic activity in Italy begin?\n"
          "https://www.polleverywhere.com/multiple_choice_polls/"
          "m1ADsg34UlgB21yiSUi3e?state=opened&flow=Default&onscreen=persist")


def s13_poll(prs):
    _poll(prs, 13, "Early History", "Where do Italians originate from?",
          "s13_image3.png",
          PN + "Where do Italians originate from?\n"
          "https://www.polleverywhere.com/multiple_choice_polls/"
          "UQ9aWvVS9TcF7EM2lW5WU?state=opened&flow=Default&onscreen=persist")


def s20_poll(prs):
    _poll(prs, 20, ROME, "How long did the Roman Empire last?",
          "s20_image8.png",
          PN + "How long did the Roman Empire last?\n"
          "https://www.polleverywhere.com/multiple_choice_polls/"
          "OSyDsfceZZWsnF4s4ouOZ?state=opened&flow=Default&onscreen=persist")


# ==========================================================================
# Original-deck image cloning  (photos/maps/charts rebuilt as native charts
# would require source data we don't have, so genuine images are preserved).
# ==========================================================================

_ORIG = Presentation(HERE / "Class 1.pptx")


def _orig_pics(idx):
    """(left, top, width, height, blob) for every picture on original slide idx."""
    out = []
    for sh in _ORIG.slides[idx - 1].shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                blob = sh.image.blob
            except Exception:
                continue
            out.append((sh.left, sh.top, sh.width, sh.height, blob))
    return out


def _place_pics(slide, pics, x0, y0, x1, y1, *, style="photo"):
    """Map the pics' bounding box into [x0,y0,x1,y1] (EMU), preserving the
    original relative arrangement and each image's aspect ratio."""
    if not pics:
        return
    L = min(p[0] for p in pics); T = min(p[1] for p in pics)
    Rr = max(p[0] + p[2] for p in pics); B = max(p[1] + p[3] for p in pics)
    bw, bh = max(1, Rr - L), max(1, B - T)
    tw, th = x1 - x0, y1 - y0
    s = min(tw / bw, th / bh)
    ox = x0 + (tw - bw * s) / 2
    oy = y0 + (th - bh * s) / 2
    for (l, t, w, h, blob) in pics:
        nl, nt = int(ox + (l - L) * s), int(oy + (t - T) * s)
        nw, nh = int(w * s), int(h * s)
        pic = slide.shapes.add_picture(BytesIO(blob), nl, nt, nw, nh)
        spPr = pic._element.spPr
        if style == "photo":
            geom = spPr.find(qn('a:prstGeom'))
            if geom is not None:
                geom.set('prst', 'roundRect')
                av = geom.find(qn('a:avLst'))
                if av is None:
                    av = ET.SubElement(geom, qn('a:avLst'))
                for gd in av.findall(qn('a:gd')):
                    av.remove(gd)
                gd = ET.SubElement(av, qn('a:gd'))
                gd.set('name', 'adj'); gd.set('fmla', 'val 5500')
            _shadow(spPr)
        elif style == "card":
            _shadow(spPr)
        # style == "flat": leave as-is


def _caption(slide, text, top, *, size=13, color=GRAY, italic=True, bold=False,
             left=MARGIN, width=RULE_W, align=PP_ALIGN.CENTER):
    _add_text(slide, left, top, width, Inches(0.4), text,
              size=size, italic=italic, bold=bold, color=color,
              font="Calibri", align=align)


# ---- generic layouts -------------------------------------------------------

def _image_slide(prs, page, section, title, orig_idx, *, style="photo",
                 source=None, caption=None, notes=None, region=None):
    slide = _blank_slide(prs)
    _top_bar(slide, section)
    top = Inches(1.5)
    if title:
        _action_title(slide, title)
    else:
        top = Inches(0.7)
    if caption:
        _caption(slide, caption, top, size=15, color=NAVY, bold=True)
        top = top + Inches(0.4)
    bot = Inches(6.85) if source else Inches(7.0)
    x0, y0, x1, y1 = region or (Inches(0.35), top, Inches(12.98), bot)
    _place_pics(slide, _orig_pics(orig_idx), x0, y0, x1, y1, style=style)
    if source:
        _caption(slide, source, Inches(6.86), size=12, color=GRAY)
    _footer(slide, page)
    if notes:
        _set_notes(slide, notes)
    return slide


def _content_clone(prs, page, section, title, items, orig_idx, *,
                   style="photo", source=None, text_w=Inches(6.8),
                   img_style=None, notes=None):
    slide = _blank_slide(prs)
    _top_bar(slide, section)
    _action_title(slide, title)
    _hbullets(slide, items, left=MARGIN, top=Inches(1.9),
              width=text_w, height=Inches(5.0))
    rx0 = MARGIN + text_w + Inches(0.25)
    _place_pics(slide, _orig_pics(orig_idx), rx0, Inches(1.65),
                Inches(13.0), Inches(6.8) if source else Inches(6.95),
                style=img_style or style)
    if source:
        _caption(slide, source, Inches(6.85), size=11, color=GRAY,
                 left=rx0, width=Inches(13.0) - rx0)
    _footer(slide, page)
    if notes:
        _set_notes(slide, notes)
    return slide


def _chart_slide(prs, page, section, title, orig_idx, *, source=None,
                 bullets=None, notes=None):
    slide = _blank_slide(prs)
    _top_bar(slide, section)
    _action_title(slide, title)
    top = Inches(1.75)
    if bullets:
        _hbullets(slide, bullets, left=MARGIN, top=top, width=RULE_W,
                  height=Inches(1.2), anchor=MSO_ANCHOR.TOP)
        top = Inches(1.75) + Inches(0.55) * len(bullets)
    bot = Inches(6.8) if source else Inches(6.95)
    _place_pics(slide, _orig_pics(orig_idx), Inches(1.2), top,
                Inches(12.13), bot, style="card")
    if source:
        _caption(slide, source, Inches(6.85), size=11, color=GRAY)
    _footer(slide, page)
    if notes:
        _set_notes(slide, notes)
    return slide


def _paper_title(prs, page, section, title, authors, note=None):
    slide = _blank_slide(prs)
    _top_bar(slide, section)
    _add_text(slide, MARGIN, Inches(2.05), RULE_W, Inches(0.5),
              "Featured research", size=16, bold=True, color=GOLD,
              font="Calibri", align=PP_ALIGN.CENTER)
    _add_text(slide, MARGIN, Inches(2.65), RULE_W, Inches(1.4), title,
              size=34, bold=True, color=NAVY, font="Calibri",
              align=PP_ALIGN.CENTER)
    aw = Inches(3.2)
    _add_rect(slide, (SLIDE_W - aw) // 2, Inches(4.35), aw, Inches(0.05), GOLD)
    _add_text(slide, MARGIN, Inches(4.65), RULE_W, Inches(0.6), authors,
              size=22, color=GRAY, font="Calibri", align=PP_ALIGN.CENTER)
    if note:
        _add_text(slide, MARGIN, Inches(5.25), RULE_W, Inches(0.5), note,
                  size=15, italic=True, color=GRAY, font="Calibri",
                  align=PP_ALIGN.CENTER)
    _footer(slide, page)
    return slide


def _blank(prs):
    _blank_slide(prs)


# ==========================================================================
# Section tags
# ==========================================================================
ROME_ = "The Roman Empire"
DARK = "The Dark Ages"
COMM = "Communes"
RENA = "The Renaissance"
DECL = "Decline & Napoleon"
UNIF = "Italian Unification"
FASC = "Fascism & WWII"
POST = "Post-WWII & Marshall Plan"
TODAY = "The Italian Economy Today"


# ==========================================================================
# Roman Empire (21-25)
# ==========================================================================

def s21(prs):
    _content(prs, 21, ROME_, "The Roman Empire", [
        (0, "Caesar’s adopted son Octavian (“Caesar Augustus”) became emperor after beating Mark Antony & Cleopatra; reigned 27 BC – 14 AD"),
        (0, "The Empire flourished over a 200-year “Pax Romana”"),
        (1, "Literature (Vergil, Horace, Ovid…)"),
        (1, "Architecture (Colosseum, Pantheon…)"),
    ])


def s22(prs):
    _image_slide(prs, 22, ROME_, "Rome’s Enduring Architecture", 22,
                 caption="Pantheon and Colosseum, Rome")


def s23(prs):
    _content(prs, 23, ROME_, "The Roman Empire", [
        (0, "Caesar’s adopted son Octavian (“Caesar Augustus”) became emperor after beating Mark Antony & Cleopatra; reigned 27 BC – 14 AD"),
        (0, "The Empire flourished over a 200-year “Pax Romana”"),
        (1, "Literature (Vergil, Horace, Ovid…); architecture (Colosseum, Pantheon…)"),
        (1, "Relatively few wars (but conquest of Britain, 43–87 AD)"),
        (1, "Trade flourished – especially silk and spices from the East"),
    ])


def s24(prs):
    _content(prs, 24, ROME_, "Decline of the Roman Empire", [
        (0, "“Five Good Emperors” until 180 AD"),
        (0, "But excessive centralization → bad emperors had immediate bad consequences"),
        (0, "Decadence and mismanagement after 200 AD"),
        (0, "Regular conflict over succession"),
    ])


def s25(prs):
    _content_clone(prs, 25, ROME_, "End of the Roman Empire", [
        (0, "Split into East and West in 395 AD"),
        (1, "West (incl. Rome) suffered frequent invasions by Germanic tribes"),
        (1, "Italy conquered by the barbarian Odoacer in 476 AD – deposed the last Western Roman emperor"),
        (1, "East continued as the Byzantine Empire through the Middle Ages"),
    ], 25, source="Coin of Odoacer, 477 AD", text_w=Inches(8.2))


# ==========================================================================
# Dark Ages (26-27)
# ==========================================================================

def s26_div(prs):
    _divider(prs, 26, 2)


def s27(prs):
    _content(prs, 27, DARK, "The Dark Ages", [
        (0, "Economic, intellectual and cultural decline across the former Western Roman Empire"),
        (0, "Constant geographic re-organization and frequent occupation"),
        (1, "Muslim rule in the South, later taken over by the Normans"),
        (1, "North integrated into the Holy Roman Empire"),
        (0, "Medieval castles all over Northern Italy"),
        (0, "Italian peninsula scattered into many small states"),
    ])


# ==========================================================================
# Communes (28-36)
# ==========================================================================

def s28_div(prs):
    _divider(prs, 28, 3)


def s29(prs):
    _content(prs, 29, COMM, "Formation of City States in the North", [
        (0, "Commercial Revolution in the 11th century"),
        (1, "Trade with the East slowly recovered (spices, silk)"),
        (0, "Cities gained economic importance"),
        (0, "Lombard League opposes Holy Roman Emperor (Frederick Barbarossa)"),
        (1, "Cities win the Battle of Legnano in 1176"),
        (1, "Treaty: cities stay loyal to the Empire but keep local jurisdiction"),
    ])


def s30(prs):
    _content(prs, 30, COMM, "Long-Run Effects of Commune Experience", [
        (0, "Independent cities had large autonomy in"),
        (1, "Regulation of economic activities"),
        (1, "Taxation"),
        (1, "Judicial power"),
        (0, "This fostered “social capital”"),
    ])


def s32(prs):
    _paper_title(prs, 32, COMM, "“Long-Term Persistence”",
                 "Luigi Guiso  ·  Paola Sapienza  ·  Luigi Zingales")


def s33(prs):
    _content(prs, 33, COMM, "Paper on Persistence of Social Capital", [
        (0, "Three proxies for social capital today"),
        (1, "Number of nonprofit organizations per capita"),
        (1, "Existence of an organ-donation organization"),
        (1, "Frequency of cheating on a national school exam"),
        (0, "Main finding: cities independent around 1200 AD have much higher social capital today"),
    ])


def s34(prs):
    _chart_slide(prs, 34, COMM, "Free City Experience and Non-Profit Organizations",
                 34, source="Mean number of non-profit organizations: 6.39")


def s35(prs):
    _chart_slide(prs, 35, COMM, "Free City Experience and Other Proxies",
                 35, source="Mean organ-donation orgs: 0.44   ·   Mean cheating in math: 2.04")


def s36(prs):
    _content(prs, 36, COMM, "Implications for Northern vs. Southern Italy", [
        (0, "Italy’s South has less “social capital” today"),
        (1, "The Normans held the South (11–13C) and blocked independent cities"),
        (0, "Southern Italy is much less developed today"),
        (1, "The gap is argued to reflect persistent differences in social capital"),
    ])


# ==========================================================================
# Renaissance (37-40)
# ==========================================================================

def s37_div(prs):
    _divider(prs, 37, 4)


def s38(prs):
    _content(prs, 38, RENA, "Renaissance (Rebirth)", [
        (0, "Independent city states prospered in the 13–16C → economic rebirth"),
        (1, "Financial innovations (banking, accounting, tradeable bonds)"),
        (1, "Shipbuilding & large fleets – new trade routes; discovery of America"),
        (0, "Wealth plus political freedom boosted art and science → cultural rebirth"),
    ])


def s39(prs):
    _image_slide(prs, 39, RENA, "Leonardo da Vinci (1452–1519)", 39)


def s40(prs):
    _image_slide(prs, 40, RENA, "Renaissance Architecture", 40,
                 caption="Milan Cathedral  ·  Marciana Library (Venice)")


# ==========================================================================
# Decline & Napoleon (41-50)
# ==========================================================================

def s41_div(prs):
    _divider(prs, 41, 5)


def s42(prs):
    _content(prs, 42, DECL, "Decline over the Period 1600–1800", [
        (0, "Discovery of America shifts activity to NW Europe → decline of Italian city states"),
        (0, "Southern Italy occupied by Spain"),
        (1, "Heavy taxes to finance Spanish wars"),
        (0, "Italy impoverished and fell behind the rest of Europe"),
    ])


def s44(prs):
    _content_clone(prs, 44, DECL, "Napoleon", [
        (0, "Spain lost its grip on Italy after 1700"),
        (1, "Last Habsburg king of Spain died without an heir"),
        (1, "Italy scattered into small states"),
        (0, "Napoleon conquered most of Italy in 1796–1809"),
        (1, "Imposed administrative and legal reforms (Code Civil)"),
        (1, "Established a sovereign Italian state in the North-East"),
    ], 44, text_w=Inches(8.3), img_style="flat")


def s45(prs):
    _content(prs, 45, DECL, "What Is the Effect of Copyright on Artistic Production?", [
        (0, "The Napoleonic setting sheds light on this important economic question"),
        (0, "Causal effect is hard to establish: copyright laws are endogenous"),
        (1, "Productive artists may lobby for protection to raise their income"),
        (0, "Yet U.S. law presumes a causal connection"),
        (1, "“The primary purpose of copyright law is to foster the creation and dissemination of intellectual works” (U.S. Copyright Laws, 1961)"),
    ])


def s46(prs):
    _paper_title(prs, 46, DECL, "“Copyright and Creativity: Evidence from Italian Operas”",
                 "Michela Giorcelli (UCLA)  ·  Petra Moser (NYU Stern)")


def s47(prs):
    _content_clone(prs, 47, DECL, "The Effect of Copyright on Creativity", [
        (0, "The Italian context helps us understand the effect of copyright on artistic production"),
        (0, "Teatro alla Scala, Milan – in the early 1800s and today"),
    ], 47, text_w=Inches(6.4))


def s48(prs):
    _content(prs, 48, DECL, "The Paper", [
        (0, "What is the causal effect of copyright on creativity?"),
        (0, "Data on 2,598 newly created operas"),
        (1, "Premiered across Italy between 1770 and 1900"),
        (1, "Measures of historical popularity and long-run durability"),
        (0, "Uses the “exogenous” imposition of copyright laws"),
        (1, "Napoleon annexes Lombardy & Venetia first: French laws (incl. copyright) in 1801"),
        (1, "Rest of Italy annexed in 1805: no copyright – a clean comparison"),
    ])


def s49(prs):
    _image_slide(prs, 49, DECL, "Napoleonic Italy: Who Got Copyright First", 49,
                 style="flat", source="Source: www.age-of-the-sage.org")


def s50(prs):
    _chart_slide(prs, 50, DECL, "Causal Effect of Copyright", 50, bullets=[
        (0, "Quantity: Lombardy & Venetia produced 2.2× more operas per year after 1801"),
        (0, "Quality: more popular and durable works"),
    ])


# ==========================================================================
# Unification (51-58)
# ==========================================================================

def s51_div(prs):
    _divider(prs, 51, 6)


def s52(prs):
    _content_clone(prs, 52, UNIF, "The Italian Unification – Failed Attempts", [
        (0, "Italy was not a unified country for 1,300 years"),
        (0, "After Napoleon’s fall (1814): again a patchwork of states"),
        (0, "Revolts against foreign rulers, aiming for a Republic"),
        (0, "Central figure: Giuseppe Garibaldi"),
        (1, "1834 uprising → sentenced to death; fled to South America"),
        (1, "Returned 1848, led the second war of independence; lost → exile in New York"),
    ], 52, text_w=Inches(9.0), img_style="photo")


def s53(prs):
    _content(prs, 53, UNIF, "The Italian Unification – Final Success", [
        (0, "Italian patriots had learned their lessons"),
        (1, "Small separate states were too weak against Austria and France"),
        (1, "No broad popular support for a Republic"),
        (0, "New approach in 1860"),
        (1, "First expel the Austrians with French help (gave Nice & Savoy to France)"),
        (1, "Back a strong Monarchy under Piedmont’s King Victor Emmanuel II"),
        (0, "Garibaldi conquered Sicily, then Naples – and handed power to the King"),
    ])


def s54(prs):
    _image_slide(prs, 54, UNIF, "Garibaldi and the Unification", 54,
                 source="Sicilian peasants welcome Garibaldi, 1860  ·  Garibaldi meets Victor Emmanuel II")


def s55(prs):
    _content_clone(prs, 55, UNIF, "Italy After the Unification", [
        (0, "Agrarian country, scarcely industrialized"),
        (0, "High illiteracy, prevalence of dialects"),
        (0, "Lack of efficient infrastructure across the peninsula"),
    ], 55, text_w=Inches(7.2),
        source="Italian peasant in the 1800s (painting by Giovanni Fattori)")


def s56(prs):
    _content_clone(prs, 56, UNIF, "Italy After the Unification", [
        (0, "Huge heterogeneity"),
        (0, "Cultural and linguistic differences"),
        (1, "Italian-speakers in 1861: about 3%"),
        (0, "Economic differences"),
        (1, "North: market economy (much richer)"),
        (1, "South: protectionist, poor, little entrepreneurial class"),
        (1, "South collapses after 1870 → the “Southern Question” and mass emigration"),
    ], 56, text_w=Inches(8.6), img_style="flat")


def s57(prs):
    _chart_slide(prs, 57, UNIF, "GDP per Capita, 1800–1945", 57,
                 source="Source: Maddison Project Database (2023 release), international-$ at 2011 prices")


def s58(prs):
    _chart_slide(prs, 58, UNIF, "Comparison, 1800–1945", 58,
                 source="Source: Maddison Project Database (2023 release), international-$ at 2011 prices")


# ==========================================================================
# Fascism & WWII (59-62)
# ==========================================================================

def s59_div(prs):
    _divider(prs, 59, 7)


def s60(prs):
    _content_clone(prs, 60, FASC, "The Rise of Fascism in Italy", [
        (0, "Early 1920s: mass strikes, worker demonstrations, high unemployment"),
        (0, "Frequent conflict between left- and right-wing militias"),
        (0, "The middle class demanded order → Mussolini promised it"),
        (1, "Oct 1922: the “March on Rome” with 30,000 Fascists"),
        (1, "The king hands power to Mussolini’s Fascist Party"),
        (1, "Mussolini turns democracy into dictatorship, builds a personality cult"),
    ], 60, text_w=Inches(9.2), img_style="photo")


def s61(prs):
    _content_clone(prs, 61, FASC, "Fascism in Italy and WWII", [
        (0, "Mussolini backed Franco in the Spanish Civil War"),
        (0, "May 1939: “Pact of Steel” with Nazi Germany"),
        (0, "But Italy was militarily weak and drained by 1943"),
        (0, "Mussolini ousted and arrested by order of the king"),
        (1, "Freed by the Germans; formed a puppet state in the North"),
        (1, "Civil war of partisans vs. Fascists and German troops"),
        (1, "German and Fascist forces surrender on 2 May"),
    ], 61, text_w=Inches(9.2), img_style="photo")


def s62_div(prs):
    _divider(prs, 62, 7)


# ==========================================================================
# Post-WWII & Marshall Plan (63-81)
# ==========================================================================

def s63_poll(prs):
    _poll(prs, 63, POST, "The Main Economic Challenge After WWII?", "s063_0.png",
          PN + "What was the main economic challenge in Italy after WWII?\n"
          "https://www.polleverywhere.com/multiple_choice_polls/"
          "7jlqfvfgQLMCXsGTvjgYo?state=opened&flow=Default&onscreen=persist")


def s64(prs):
    _content_clone(prs, 64, POST, "Capital and Management in Post-WWII Italy", [
        (0, "10% of physical capital damaged by WWII"),
        (0, "Antiquated machines; poor managerial practices"),
        (1, "Plants poorly organized; little maintenance; old-fashioned marketing (BLS, 1949)"),
        (1, "“Lack of management was a more severe problem than war damages” (Silberman et al., 1996)"),
    ], 64, text_w=Inches(8.4),
        source="Reconstruction of the Ariccia Bridge, Rome, 1946–1948")


def s65(prs):
    _content_clone(prs, 65, POST, "The Marshall Plan", [
        (0, "Massive aid to help Western Europe recover from WWII"),
        (1, "Financial aid and in-kind subsidies (1948–1951)"),
        (1, "Aid of $130 billion (2010 USD) – 5% of US GDP in 1948"),
    ], 65, text_w=Inches(8.6), img_style="flat")


def s66(prs):
    _image_slide(prs, 66, POST, "The Marshall Plan in the Public Eye", 66,
                 caption="Contemporary posters and cartoons")


def s67(prs):
    _content(prs, 67, POST, "The Marshall Plan Productivity Program in Italy", [
        (0, "Productivity Program (1952–1958)"),
        (1, "Management training trips for European managers at US firms"),
        (1, "Loans restricted to technologically advanced US machines"),
        (0, "Implementation in Italy"),
        (1, "Targeted small and medium-sized manufacturing firms"),
        (1, "Participation on a voluntary basis"),
    ])


def s68(prs):
    _content_clone(prs, 68, POST, "Management Practices Taught During Study Trips", [
        (0, "Training within Industry (TWI)"),
        (0, "Training trips"),
        (1, "Teams of 15–20 people from European countries"),
        (1, "8–12 weeks in 5–6 US firms in the same sector"),
        (1, "Seminars with US experts; working side-by-side with US managers"),
    ], 68, text_w=Inches(9.0), img_style="photo")


def s69(prs):
    _content(prs, 69, POST, "Management Practices Taught During Study Trips", [
        (0, "Factory operations"),
        (1, "Regular machine maintenance; workplace safety"),
        (0, "Production planning"),
        (1, "Sales and order control"),
        (0, "HR training and management"),
        (1, "Training & supervision; continuous improvement of methods"),
        (0, "Marketing"),
        (1, "Market research, branding, design; modern advertising & distribution"),
    ])


def s70(prs):
    _content_clone(prs, 70, POST, "Firm Financing and Purchase of US Machines", [
        (0, "Subsidized loans"),
        (1, "Buy modern US machines at low interest rates"),
        (0, "4–7 week study trips for Italian engineers"),
        (1, "Transfer of the know-how to use the machines"),
        (1, "3-year technical assistance in Italy"),
    ], 70, text_w=Inches(8.6),
        source="Fiat factory using new US Marshall Plan machinery")


def s71(prs):
    _content(prs, 71, POST, "Did the Marshall Plan Raise Productivity in Italy?", [
        (0, "Challenge: firm participation was voluntary"),
        (0, "Remember the “endogeneity” problem?"),
        (0, "Firms that signed up were probably already different"),
        (1, "More motivated managers, keener on modern equipment"),
        (1, "The “treatment” group can’t be directly compared to the “control” group"),
    ])


def s72(prs):
    _paper_title(prs, 72, POST,
                 "“The Long-Term Effects of Management and Technology Transfers”",
                 "Michela Giorcelli (UCLA)")


def s73(prs):
    _content(prs, 73, POST, "The Paper", [
        (0, "What was the impact of the program on firm performance?"),
        (0, "Data on 6,065 Italian firms eligible to participate"),
        (1, "Balance sheets from 5 years before to 15 years after"),
        (0, "Method to identify the causal effect"),
        (1, "Uses unexpected (“exogenous”) US budget cuts"),
        (1, "Treatment: firms that still got funding/training"),
        (1, "Control: eligible applicants excluded after the cut"),
    ])


def s74(prs):
    _image_slide(prs, 74, POST, "SME Manufacturing Firms in 5 Pilot Regions", 74,
                 style="flat", source="Pilot regions vs. rest of Italy")


def s75(prs):
    _image_slide(prs, 75, POST, "The Unexpected Budget Cut: 5 Treated Provinces", 75,
                 style="flat", source="Pilot regions  ·  treated provinces within them")


def s76(prs):
    _chart_slide(prs, 76, POST, "Impact of Managerial Training (i): Higher Survival",
                 76, source="Treatment vs. comparison provinces")


def s77(prs):
    _chart_slide(prs, 77, POST, "Impact of Managerial Training (ii): Sales & Productivity",
                 77, source="Log TFPR, years after intervention")


def s78(prs):
    _chart_slide(prs, 78, POST, "Technology Alone Had Only a Small Effect", 78,
                 source="Impact of technology transfer, years after US intervention")


def s79(prs):
    _content(prs, 79, POST, "Mechanisms Behind the Results", [
        (0, "What changed in the management of TWI firms?"),
        (0, "90% adopted US management practices – still in use 15 years later"),
        (0, "Better organization: higher manager-to-worker ratio"),
        (0, "Higher loans, investment, and ROA"),
    ])


def s80(prs):
    _chart_slide(prs, 80, POST, "Detailed Data on Implementation", 80,
                 source="Share of firms that adopted US managerial practices, by year")


def s81(prs):
    _content(prs, 81, POST, "Main Take-Aways from the Paper", [
        (0, "Improving managerial practices had a major impact on sales and productivity"),
        (0, "Underlines the importance of business training"),
        (0, "Post-WWII Italy resembles many developing countries today"),
    ])


# ==========================================================================
# The Italian Economy Today (82-89)
# ==========================================================================

def s82_div(prs):
    _divider(prs, 82, 8)


def s83(prs):
    _chart_slide(prs, 83, TODAY, "GDP after WWII", 83, bullets=[
        (0, "Unprecedented growth after WWII, then stagnation since the 2000s"),
    ], source="Source: Maddison Project Database (2023 release), international-$ at 2011 prices")


def s84(prs):
    _chart_slide(prs, 84, TODAY, "Comparison: Italy, Germany, United States", 84,
                 source="Source: Maddison Project Database (2023 release), international-$ at 2011 prices")


def s85(prs):
    _chart_slide(prs, 85, TODAY, "Total Factor Productivity", 85, bullets=[
        (0, "Italy’s productivity has been stagnant since the mid-1990s – slow IT adoption? misallocation?"),
    ], source="Source: Calligaris et al. (2018), The Productivity Puzzle and Misallocation")


def s86(prs):
    _chart_slide(prs, 86, TODAY, "Huge Differences North–South", 86, bullets=[
        (0, "Annual income per capita, 2014"),
    ])


def s87(prs):
    _content_clone(prs, 87, TODAY, "Brain Drain and Out-Migration", [
        (0, "Among the EU countries sending the most workers abroad"),
        (0, "2024: ~156,000 citizens emigrated – a 25-year high; graduates especially"),
        (0, "New research: graduate emigration → fewer start-ups created back home"),
    ], 87, text_w=Inches(6.6), img_style="flat",
        source="ISTAT (2024)  ·  Anelli, Basso, Ippedico & Peri (AEJ: Applied, 2023)")


def s88(prs):
    _content_clone(prs, 88, TODAY, "Corruption", [
        (0, "Ranked 52 of 180 on the 2024 Corruption Perception Index (score 54/100)"),
        (0, "19th of 27 in the EU – dropped 10 places vs. 2023"),
    ], 88, text_w=Inches(6.4), img_style="flat",
        source="Source: Transparency International, Corruption Perception Index 2024")


def s89(prs):
    _content_clone(prs, 89, TODAY, "Strong Brands!", [
        (0, "Prominent sectors"),
        (1, "Food"),
        (1, "Tourism"),
        (1, "Fashion"),
        (1, "Luxury cars"),
    ], 89, text_w=Inches(4.6), img_style="flat")


# ==========================================================================
# Quiz / close / topics (90-93)
# ==========================================================================

def s90(prs):
    slide = _blank_slide(prs)
    _top_bar(slide, "Quiz")
    _action_title(slide, "Quiz: “The Good, the Bad, and the Ugly”")
    _hbullets(slide, [
        (0, "Give two examples of the Italian economy’s…"),
        (1, "“good” attributes"),
        (1, "“bad” attributes"),
        (1, "“ugly” attributes"),
    ], top=Inches(2.0), height=Inches(4.2))
    _footer(slide, 90)


def s91(prs):
    slide = _blank_slide(prs)
    _add_text(slide, MARGIN, Inches(3.0), RULE_W, Inches(1.3), "Thank You!",
              size=60, bold=True, color=NAVY, font="Calibri", align=PP_ALIGN.CENTER)
    aw = Inches(3.0)
    _add_rect(slide, (SLIDE_W - aw) // 2, Inches(4.3), aw, Inches(0.06), GOLD)
    _add_rect(slide, 0, Inches(7.15), SLIDE_W, Inches(0.02), RULE)
    _add_rect(slide, MARGIN, Inches(7.135), GOLD_W, Inches(0.05), GOLD)


def _section_cover(prs, page, title):
    slide = _blank_slide(prs)
    _top_bar(slide, title)
    _add_text(slide, MARGIN, Inches(3.0), RULE_W, Inches(1.0), title,
              size=44, bold=True, color=NAVY, font="Calibri", align=PP_ALIGN.CENTER)
    aw = Inches(3.0)
    _add_rect(slide, (SLIDE_W - aw) // 2, Inches(4.2), aw, Inches(0.05), GOLD)
    _footer(slide, page)


def s92(prs):
    _section_cover(prs, 92, "Presentation Topics")


def s93(prs):
    _image_slide(prs, 93, "Presentation Topics", "Presentation Topics", 93,
                 style="flat")


# ==========================================================================
# Backup slides (94-120)
# ==========================================================================
BK = "Backup"


def s94(prs):
    _section_cover(prs, 94, "Backup Slides")


def s95(prs):
    _image_slide(prs, 95, BK, "Trip Schedule (Sep 6–12)", 95, style="flat")


def s96(prs):
    _image_slide(prs, 96, BK, "Villa Campari", 96)


def s97(prs):
    _image_slide(prs, 97, BK, "Manifattura Automobili Torino (MAT)", 97)


def s98(prs):
    _image_slide(prs, 98, BK, "Matera – the Sassi", 98)


def s99(prs):
    _image_slide(prs, 99, BK, "Valcamonica Rock Drawings", 99, style="flat")


def s100(prs):
    _image_slide(prs, 100, BK, "“The Astronauts” – Valcamonica", 100, style="flat")


def s101(prs):
    _image_slide(prs, 101, BK, "Etruscan & Greek Territories", 101, style="flat")


def s102(prs):
    _image_slide(prs, 102, BK, "Etruscan Art", 102)


def s103(prs):
    _image_slide(prs, 103, BK, "San Marino", 103)


def s106(prs):
    _image_slide(prs, 106, BK, "Member Cities of the Lombard Leagues", 106, style="flat")


def s107(prs):
    _image_slide(prs, 107, BK, "Librettos and Opera Scores (1800s)", 107)


def s108(prs):
    _content_clone(prs, 108, BK, "Measuring Opera Quality", [
        (0, "Historical popularity: listed in the Annals of Opera (Loewenberg, 1978)"),
        (0, "Performed at the Met in NYC between 1900 and 2014"),
        (0, "A complete recording available on Amazon today"),
    ], 108, text_w=Inches(6.6), img_style="flat",
        source="La Traviata (Verdi, 1853)  ·  Lucia di Lammermoor (Donizetti, 1835)")


def s120_poll(prs):
    _poll(prs, 120, BK, "Italian-Speakers in Italy in 1861?", "s120_0.png",
          PN + "What was the percentage of people in Italy speaking Italian in 1861?\n"
          "https://www.polleverywhere.com/multiple_choice_polls/"
          "XVz6lRn9EFqmrk76BjE4T?state=opened&flow=Default&onscreen=persist")


# History-of-Italy map screenshots (104,105,109-119) + blanks (121-123)
def _bk_map(page, idx):
    return lambda prs: _image_slide(prs, page, BK, None, idx, style="flat")


# ==========================================================================
# Poll builders for the Communes / Napoleon sections
# ==========================================================================

def s31_poll(prs):
    _poll(prs, 31, COMM, "What Is “Social Capital”?", "s031_0.png",
          PN + "What is \"Social Capital\"?\n"
          "https://www.polleverywhere.com/multiple_choice_polls/"
          "0Vakih1wuXkrhOd1qaRi6?state=opened&flow=Default&onscreen=persist")


def s43_poll(prs):
    _poll(prs, 43, DECL, "Napoleon’s Occupation of Italy…", "s043_0.png",
          PN + "Napoleon's occupation...\n"
          "https://www.polleverywhere.com/multiple_choice_polls/"
          "d49GVIcC2Dw0TNwAOAQUr?state=opened&flow=Default&onscreen=persist")


# ==========================================================================
# Newer papers added on 2026-07-18 (paper-writing-agent review)
# ==========================================================================

def paper_bianchi_title(prs, page=82):
    _paper_title(prs, page, POST,
                 "“Reconstruction Aid, Public Infrastructure, and Economic Development”",
                 "Nicola Bianchi  ·  Michela Giorcelli",
                 note="Journal of Economic History, 2023")


def paper_bianchi_find(prs, page=83):
    _content(prs, page, POST, "The Marshall Plan’s “Hardware”: Infrastructure", [
        (0, "The Marshall Plan also rebuilt roads, bridges, power, and water – not just firms"),
        (0, "Provinces that received more reconstruction grants grew faster for decades"),
        (0, "Infrastructure (“hardware”) and management training (“software”) reinforced each other"),
        (0, "Takeaway: rebuilding physical capital and know-how together drives lasting growth"),
    ])


def paper_pz_title(prs, page=88):
    _paper_title(prs, page, TODAY,
                 "“Diagnosing the Italian Disease”",
                 "Bruno Pellegrino  ·  Luigi Zingales",
                 note="NBER Working Paper 23964")


def paper_pz_find(prs, page=89):
    _content(prs, page, TODAY, "Why Italian Productivity Stalled", [
        (0, "Italy’s productivity stopped growing in the mid-1990s – just as IT took off"),
        (0, "The cause was largely internal: family firms promote on loyalty, not merit"),
        (0, "They bought computers but couldn’t use them well – managers lacked the skills"),
        (0, "Takeaway: without meritocratic management, new technology doesn’t pay off"),
    ])


# ==========================================================================
# Build order
# ==========================================================================
BUILD_ORDER = [
    s01_cover, s02_title, s03, s04, s05, s06, s07, s08, s09, s10_roadmap,
    s11_poll, s12_divider, s13_poll, s14, s15, s16, s17_divider, s18, s19, s20_poll,
    s21, s22, s23, s24, s25, s26_div, s27, s28_div, s29, s30,
    s31_poll, s32, s33, s34, s35, s36, s37_div, s38, s39, s40,
    s41_div, s42, s43_poll, s44, s45, s46, s47, s48, s49, s50,
    s51_div, s52, s53, s54, s55, s56, s57, s58, s59_div, s60,
    s61, s62_div, s63_poll, s64, s65, s66, s67, s68, s69, s70,
    s71, s72, s73, s74, s75, s76, s77, s78, s79, s80,
    s81, paper_bianchi_title, paper_bianchi_find,
    s82_div, s83, s84, s85, paper_pz_title, paper_pz_find,
    s86, s87, s88, s89, s90,
    s91, s92, s93, s94, s95, s96, s97, s98, s99, s100,
    s101, s102, s103, _bk_map(104, 104), _bk_map(105, 105), s106, s107, s108,
    _bk_map(109, 109), _bk_map(110, 110), _bk_map(111, 111), _bk_map(112, 112),
    _bk_map(113, 113), _bk_map(114, 114), _bk_map(115, 115), _bk_map(116, 116),
    _bk_map(117, 117), _bk_map(118, 118), _bk_map(119, 119), s120_poll,
    _blank, _blank, _blank,
]

# Poll embed GUIDs and internal links, keyed to BUILD functions so the
# splice/wire scripts stay correct no matter where slides are inserted.
POLL_GUIDS = {
    s11_poll: "9ac8ab11-17d4-4f8b-ad71-6ec62cf9d0d7",
    s13_poll: "04b94b40-31c3-4b7c-9d15-6c2cfb785dbe",
    s20_poll: "d31d8fc9-0ae4-4dd9-bcb7-12494efd565e",
    s31_poll: "5d3153f0-a344-4011-8a54-16cb04a9bffa",
    s43_poll: "debcd7cc-5e9f-4b3f-b122-b4a0b5db01df",
    s63_poll: "8fa44402-3d72-47bf-922e-7724c3dabf1c",
    s120_poll: "2cee3236-b254-4eb9-8edb-1ca5f35ef05a",
}
LINK_SPECS = [
    (s08, "topics list", s93),
    (s14, "Sassi", s98),
    (s14, "The Astronauts", s99),
    (s15, "metal trade", s101),
    (s15, "Advanced art", s102),
    (s29, "Lombard League", s106),
    (s04, "Schedule", s95),
]


def main():
    import json
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    for fn in BUILD_ORDER:
        fn(prs)
    prs.save(OUT)
    n = len(prs.slides._sldIdLst)

    def disp(fn):
        return BUILD_ORDER.index(fn) + 1
    manifest = {
        "polls": {guid: disp(fn) for fn, guid in POLL_GUIDS.items()},
        "links": [[disp(s), needle, disp(t)] for s, needle, t in LINK_SPECS],
    }
    (HERE / "_manifest.json").write_text(json.dumps(manifest, indent=2))
    print(f"Wrote {OUT}  ({n} slides); manifest updated")


if __name__ == "__main__":
    main()
